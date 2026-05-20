"""
Full 43-slide build of Лекции 9 «AI в авиакосмической отрасли и оборонном комплексе».

Source-of-truth: deck.yaml v1 + chapter v3 (status=finalized, ~17k слов, 28-term глоссарий,
OODA keystone §0.2, 104 source) + slides/*.md (43 файла с readable speaker notes).

Issue #118 · downstream-каскад от chapter (book-first).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml ≈ 75 мин.

Lec-N-1 pattern compliance: match lec-07 (cover + lecture-map + 6 section dividers +
dedicated Q&A; top progress bar только на dividers + cover).

Build via: python3 build_lec09.py — generates lec-09.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
LIGHT_TINT = RGBColor(0xE6, 0xEE, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
RED_WARN = RGBColor(0xC0, 0x39, 0x2B)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered" / "assets"
OUT = ROOT / "rendered" / "lec-09.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


# ========== Helpers ==========

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def hr_line(slide, x, y, w, color=LIGHT, weight=1.0):
    shp = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    shp.line.color.rgb = color
    shp.line.width = Pt(weight)
    return shp


def add_arrow(slide, x, y, w, h, *, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w), height=Inches(h))
    elif w:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w))
    else:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


# ========== Slide-level helpers ==========

def add_progress_bar(slide, current_section, *, y=6.95):
    """6-card progress bar для section dividers и cover.
    current_section=-1 → no card highlighted (e.g., cover slide)."""
    sections = ["Открытие", "Sense", "Decide", "Act", "Граница", "Сборка"]
    n = len(sections)
    bar_x = 0.5
    bar_w = 12.33
    gap = 0.08
    card_w = (bar_w - gap * (n - 1)) / n
    card_h = 0.32
    for i, name in enumerate(sections):
        cx = bar_x + i * (card_w + gap)
        is_current = (i == current_section and current_section >= 0)
        if is_current:
            ocean_box(slide, cx, y, card_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=0.0, radius_pt=6)
            text_box(slide, cx, y+0.04, card_w, card_h-0.06,
                     f"{i}. {name}", size=10, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            ocean_box(slide, cx, y, card_w, card_h,
                      fill=SURFACE, stroke=LIGHT, stroke_pt=0.5, radius_pt=6)
            text_box(slide, cx, y+0.04, card_w, card_h-0.06,
                     f"{i}. {name}", size=10, bold=False, color=LIGHT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text):
    """Standard footer для caption / source 12pt italic."""
    text_box(slide, 0.5, 7.05, 12.33, 0.30, text,
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)


def add_assertion_title(slide, text, *, size=26, y=0.4):
    """Assertion as page title — 26pt bold deep, left aligned."""
    text_box(slide, 0.6, y, 12.13, 1.2, text,
             size=size, bold=True, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


# ========== Slide builders ==========

def slide_01_hook(prs):
    """s01 — BEFORE/AFTER satellite hook"""
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "ИИ-аналитика спутниковой съёмки превращает 250 ПБ архива Maxar\n"
        "в детекции за часы — то, что раньше требовало дней.",
        size=24)

    # Left: real satellite imagery (Sentinel-2 from Wikimedia Commons)
    ocean_box(s, 0.6, 2.0, 7.5, 4.6)
    photo_p = ASSETS / "photos" / "p01-sat-imagery.jpg"
    if photo_p.exists():
        add_image(s, photo_p, 0.8, 2.2, w=7.1, h=4.0)
        # Annotated overlay: gold rectangle highlighting "change region"
        bb1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(3.7),
                                 Inches(1.5), Inches(0.9))
        bb1.fill.background()
        bb1.line.color.rgb = GOLD; bb1.line.width = Pt(2.5)
        text_box(s, 2.7, 4.65, 1.8, 0.25, "обнаруженное\nизменение", size=9,
                 italic=True, color=GOLD, bold=True, line_spacing=1.1)
        text_box(s, 0.85, 6.3, 7.0, 0.3,
                 "Sentinel-2 (ESA, Wikimedia Commons CC-BY-SA) — иллюстрация формата спутниковых данных, на которых работает Maxar Sentry.",
                 size=9, italic=True, color=LIGHT, line_spacing=1.2)

    # Right: info-card sidebar
    ocean_box(s, 8.5, 2.0, 4.3, 4.6, fill=LIGHT_TINT)
    text_box(s, 8.75, 2.2, 3.8, 0.4, "Maxar Sentry · 2025",
             size=18, bold=True, color=DEEP)
    text_box(s, 8.75, 2.7, 3.8, 0.35,
             "Платформа упреждающей разведки", size=13, italic=True, color=MID)
    hr_line(s, 8.75, 3.15, 3.7, color=LIGHT, weight=0.75)

    text_runs(s, 8.75, 3.3, 3.8, 0.7, [
        {"text": "250 ПБ", "size": 30, "bold": True, "color": GOLD},
        {"newpara": True, "text": "архив спутниковых снимков, 20+ лет",
         "size": 11, "italic": True, "color": MID},
    ])

    text_runs(s, 8.75, 4.3, 3.8, 0.6, [
        {"text": "3 сенсора", "size": 22, "bold": True, "color": DEEP},
        {"newpara": True, "text": "оптика + радар + AIS судов · перекрёстное наведение",
         "size": 11, "italic": True, "color": MID, "line_spacing": 1.2},
    ])

    text_runs(s, 8.75, 5.1, 3.8, 0.8, [
        {"text": "Часы", "size": 28, "bold": True, "color": GOLD},
        {"text": " после съёмки", "size": 14, "color": DEEP},
        {"newpara": True, "text": "не дни, как было до 2022", "size": 11, "italic": True, "color": MID},
    ])

    text_box(s, 8.75, 5.95, 3.8, 0.5, "Контракт NGA Luno A D01",
             size=11, italic=True, color=LIGHT)

    add_footer(s,
        "Источники: Defense One, июнь 2025 · BusinessWire 20250625 · Military Aerospace 2025")

    add_speaker_notes(s, _notes_s01())
    return s


def _notes_s01():
    return ("Представьте обычный понедельник мая 2026 года. Аналитик утром "
            "открывает сводку Maxar Sentry и видит около двух десятков аномалий, "
            "которые система автоматически пометила за прошедшую ночь. Среди них — "
            "пара зданий в новом порту в Африке, появившихся между съёмками с "
            "интервалом в восемь часов.\n\n"
            "Что важно в этой картинке для нас, инженеров? Во-первых, аналитика "
            "была сделана не глазами человека. Maxar Sentry обучен на архиве "
            "объёмом около 250 петабайт спутниковых снимков высокого разрешения "
            "за два десятилетия, и его конвейер постоянно сравнивает свежие "
            "снимки с базовой линией. Во-вторых, человек всё ещё в петле: он "
            "смотрит на «before/after», подтверждает или отвергает гипотезу. "
            "AI ускорил Sense, но Decide и Act остались на людях. И в-третьих — "
            "это работает в часах от съёмки до сводки, а не в днях.\n\n"
            "Этот кадр — не победа AI и не его поражение. Это маленький фрагмент "
            "длинной цепи, в которой одни звенья AI отдают людям, а в других "
            "звеньях люди сознательно оставляют контроль за собой. Именно эту "
            "цепь мы сейчас разложим как карту лекции.")


def slide_02_cover(prs):
    """s02 — Cover. Decorative «09» + title + meta."""
    s = blank(prs); set_slide_bg(s, WHITE)

    # Decorative «09» outline — huge, lower-left
    text_box(s, 0.0, 0.5, 6.5, 7.0, "09",
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD, line_spacing=0.9)

    # Meta line — top right
    text_box(s, 6.0, 1.5, 6.8, 0.5, "Лекция 9 · AI в инженерных задачах · 3 курс",
             size=18, italic=True, color=MID)

    # Title — large
    text_box(s, 6.0, 2.3, 6.8, 3.0,
             "AI в авиакосмической отрасли\nи оборонном комплексе",
             size=36, bold=True, color=DEEP, line_spacing=1.15)

    # Duration
    text_box(s, 6.0, 5.7, 6.8, 0.5, "75 минут · Q&A",
             size=18, italic=True, color=LIGHT)

    # Small radar/satellite decoration
    icon_p = ASSETS / "icons" / "radar-96.png"
    if icon_p.exists():
        add_image(s, icon_p, 11.5, 6.3, w=0.9, h=0.9)

    # P2-1: Lec-07 pattern — cover BEZ progress bar (только dividers + content где включено)

    add_speaker_notes(s,
        "Это титульный слайд лекции 9 — седьмой отраслевой в нашем курсе "
        "и одной из двух (вместе с медициной), где цена единственной ошибки "
        "модели может измеряться не деньгами, а жизнями. В отличие от медицины, "
        "где регулятор единый — FDA или EU AI Act — в аэрокосмосе и обороне "
        "инженер одновременно встроен в три рамки: гражданская сертификация "
        "авиационного ПО, оборонная политика отдельной страны и международное "
        "право. Эти три рамки иногда совпадают, иногда конкурируют, и навык "
        "находить свой путь между ними — отдельная инженерная компетенция, "
        "которую мы будем тренировать сегодня. Тон лекции — trust-but-verify: "
        "не евангелизм и не диссидентство. Мы разбираем, что именно работает, "
        "где именно ломается и где именно проходит граница, которую инженер "
        "не пересекает по решению вендора, а только по решению государства и "
        "международного сообщества.")
    return s


def slide_03_lecture_map(prs):
    """s03 — Lecture-map. 6 cards + central question."""
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Карта лекции — 4 содержательных раздела плюс граница и сборка")

    # 6 cards in 2 rows × 3 cols
    cards = [
        ("0. Открытие", "Hook + keystone OODA + glossary", ""),
        ("1. Sense", "Глаза в небе и на земле", "12 мин"),
        ("2. Decide", "От наблюдения к решению", "14 мин"),
        ("3. Act", "Автономия на платформе", "14 мин"),
        ("4. Граница", "и регулирование (LAWS, UN GGE)", "15 мин"),
        ("5. Сборка", "критерии, карьера, замыкание", "6 мин"),
    ]
    col_w = 4.0; col_h = 1.8; gap = 0.2
    start_x = 0.6
    for i, (title, desc, dur) in enumerate(cards):
        row = i // 3; col = i % 3
        x = start_x + col * (col_w + gap)
        y = 1.8 + row * (col_h + gap)
        ocean_box(s, x, y, col_w, col_h)
        text_box(s, x + 0.25, y + 0.15, col_w - 0.5, 0.5, title,
                 size=20, bold=True, color=DEEP)
        text_box(s, x + 0.25, y + 0.75, col_w - 0.5, 0.7, desc,
                 size=13, color=MID, italic=True, line_spacing=1.2)
        if dur:
            text_box(s, x + 0.25, y + col_h - 0.5, col_w - 0.5, 0.35, dur,
                     size=11, italic=True, color=GOLD, bold=True,
                     align=PP_ALIGN.RIGHT)

    # Central question banner
    ocean_box(s, 0.6, 5.7, 12.13, 1.1, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.85, 5.78, 11.63, 0.4,
             "Центральный вопрос лекции:",
             size=13, italic=True, color=DARK_GREY, bold=True)
    text_box(s, 0.85, 6.13, 11.63, 0.6,
             "Где в цепи Sense → Decide → Act AI даёт измеримое преимущество — "
             "и где он становится опасной автоматизацией?",
             size=14, italic=True, color=DEEP)

    add_speaker_notes(s,
        "Структура лекции повторяет цепь OODA с одним важным добавлением: там, "
        "где звено Act обрезано регулированием, мы посвящаем отдельный раздел. "
        "Первый раздел — Sense. Это место, где AI работает лучше всего: данных "
        "много, ground truth относительно доступна, цена ложного срабатывания "
        "терпима. Здесь разберём спутниковую аналитику Maxar и BlackSky, edge-AI "
        "на орбите, predictive maintenance Rolls-Royce и Airbus Skywise, "
        "российский слой ТЕРРА ТЕХ и СКАНЭКС — а также три провала. Второй — "
        "Decide. Это самое тонкое звено: спрос на ускорение огромен, но "
        "галлюцинации больших языковых моделей и automation bias превращаются "
        "здесь из лабораторных феноменов в решения о жизнях людей. Здесь "
        "Palantir MSS, Scale Donovan, Helsing, Anthropic-Palantir partnership — "
        "и канонический разбор Lavender как «accuracy — не та метрика». Третий — "
        "Act. Хайп далеко впереди реальности. CCA, X-62A, Saker Scout, Geran-2 — "
        "реальные кейсы; MCAS как анти-паттерн safety-critical. Четвёртый — "
        "граница и регулирование. Лестница L1-L5, UN GGE, ICRC, Maven walkout и "
        "big-tech возврат, триада HITL/HOOL/HOTL. Пятый — сборка: семь критериев, "
        "карьерный угол, чтение, замыкание. Центральный вопрос внизу — держите "
        "его в голове по ходу всей лекции.")
    return s


def slide_04_glossary(prs):
    """s04 — Glossary mini, 6 acronyms × 2 cols."""
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, "Шесть аббревиатур, без которых дальше не пройти")

    acronyms = [
        ("SAR", "Synthetic Aperture Radar",
         "Активный радар, видит сквозь облака и ночью."),
        ("ATR", "Automatic Target Recognition",
         "ML-классификатор объектов по сенсорному образу."),
        ("ISR", "Intelligence, Surveillance, Reconnaissance",
         "Зонтичный термин для разведзадач."),
        ("EW", "Electronic Warfare",
         "Радиоэлектронная борьба: подавление, обман, перехват."),
        ("LAWS", "Lethal Autonomous Weapon Systems",
         "Летальные автономные системы, предмет UN GGE."),
        ("OODA", "Observe-Orient-Decide-Act",
         "Цикл принятия решения Boyd 1976."),
    ]
    col_w = 5.95; col_h = 1.55; gap = 0.2
    for i, (acr, full, defn) in enumerate(acronyms):
        row = i // 2; col = i % 2
        x = 0.6 + col * (col_w + gap)
        y = 1.7 + row * (col_h + gap)
        ocean_box(s, x, y, col_w, col_h)
        text_box(s, x + 0.25, y + 0.15, col_w - 0.5, 0.5, acr,
                 size=24, bold=True, color=DEEP)
        text_box(s, x + 1.8, y + 0.25, col_w - 1.95, 0.4, full,
                 size=11, italic=True, color=LIGHT)
        text_box(s, x + 0.25, y + 0.75, col_w - 0.5, 0.7, defn,
                 size=13, color=MID, line_spacing=1.2)

    add_footer(s,
        "Остальные аббревиатуры (HITL, HOOL, HOTL, CCA, MCAS, IFF и др.) — "
        "расшифровка при первом упоминании по ходу лекции")

    add_speaker_notes(s,
        "В этой лекции мы пользуемся шестью ключевыми аббревиатурами на каждом "
        "шагу. Этот слайд — не для запоминания наизусть; он работает как "
        "табличка-якорь. SAR — это радар с синтезированной апертурой; он "
        "формирует изображение по принципу пересчёта эха со множества положений "
        "антенны. Это главная альтернатива оптической съёмке для всепогодного "
        "мониторинга. ATR — автоматическое распознавание целей; это "
        "ML-классификатор, который по сенсорному изображению определяет тип и "
        "идентичность объекта. Главный «декорированный» термин в военной "
        "AI-рекламе. ISR — зонтичный термин для разведзадач: спутниковая, "
        "воздушная, наземная разведка плюс долгосрочное наблюдение. EW — "
        "радиоэлектронная борьба, класс действий по подавлению, обману или "
        "перехвату сигналов противника; именно EW обычно атакует AI-системы "
        "через сенсорный вход. LAWS — летальные автономные системы оружия; "
        "термин, вокруг которого идут переговоры в UN GGE. OODA — цикл Боуда, "
        "на котором построена вся эта глава. Остальные аббревиатуры — HITL, "
        "HOOL, HOTL, CCA, MCAS, IFF, AMRAAM, V-BAT, GNSS, FedRAMP, ROE, IHL — "
        "будут расшифрованы по первому упоминанию.")
    return s


def slide_05_keystone(prs):
    """s05 — Keystone OODA chain"""
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, "Три звена цепи. ИИ входит в каждое — но по-разному.")

    # Dual-use band
    filled_rect(s, 0.6, 1.65, 12.13, 0.3, SOFT_GREY)
    text_box(s, 0.6, 1.7, 12.13, 0.25,
             "Гражданское ↔ Военное · одни модели, два контура применения",
             size=11, italic=True, color=DARK_GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 3 OODA cards
    card_w = 3.7; card_h = 3.6; gap = 0.55
    start_x = 0.6 + (12.13 - 3*card_w - 2*gap) / 2

    items = [
        ("Sense", "eye", LIGHT, [
            ("ИИ работает лучше всего", DEEP),
            ("Эталонная разметка доступна", MID),
            ("Терпимая цена ложной тревоги", MID),
        ]),
        ("Decide", "brain", MID, [
            ("ИИ как ускоритель аналитика — хорошо", DEEP),
            ("ИИ как замена аналитика — плохо", GOLD),
            ("Здесь LLM-хайп опаснее всего", MID),
        ]),
        ("Act", "plane", DEEP, [
            ("Узкие сценарии с пилотом-наблюдателем", DEEP),
            ("Полная автономия — маркетинг 2026", MID),
            ("Большинство ударов — с оператором в петле", MID),
        ]),
    ]
    for i, (name, icon_name, accent, lines) in enumerate(items):
        x = start_x + i * (card_w + gap)
        y = 2.2
        ocean_box(s, x, y, card_w, card_h)

        # Icon
        icon_p = ASSETS / "icons" / f"{icon_name}-96.png"
        if icon_p.exists():
            add_image(s, icon_p, x + (card_w - 0.9) / 2, y + 0.3, w=0.9, h=0.9)
        else:
            # fallback — circle
            ic = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x + (card_w - 0.9) / 2),
                                    Inches(y + 0.3),
                                    Inches(0.9), Inches(0.9))
            ic.fill.solid(); ic.fill.fore_color.rgb = MID
            ic.line.fill.background()

        # Title
        text_box(s, x + 0.2, y + 1.35, card_w - 0.4, 0.5, name,
                 size=28, bold=True, color=accent,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        # Lines
        for j, (line_text, color) in enumerate(lines):
            text_box(s, x + 0.3, y + 2.0 + j * 0.5, card_w - 0.6, 0.45,
                     "• " + line_text,
                     size=12, color=color, line_spacing=1.2)

        # Arrows between cards (after card 1 and 2)
        if i < 2:
            add_arrow(s, x + card_w + 0.05, y + card_h / 2 - 0.2,
                      gap - 0.1, 0.4, fill=LIGHT)

    add_footer(s,
        "Бойд, ВВС США, 1976 — изначально модель воздушного боя; "
        "теперь универсальная рамка цикла принятия решений.")

    add_speaker_notes(s,
        "Любая аэрокосмическая или оборонная задача — будь то перехват ракеты, "
        "корректировка курса по погоде, мониторинг изменений на территории — "
        "может быть разложена на одну и ту же цепочку шагов. Сначала надо "
        "увидеть: собрать сигналы из мира с сенсоров, камер, радаров, открытых "
        "источников. Это Sense. Затем эти сигналы надо проинтерпретировать и "
        "решить: классифицировать объект, оценить угрозу, выбрать вариант "
        "реакции. Это Decide. И наконец что-то сделать: дать команду оператору, "
        "развернуть систему ПВО, спустить дрон. Это Act.\n\n"
        "Эту трёхзвенную модель в 1976 году сформулировал американский военный "
        "лётчик и теоретик полковник Джон Бойд. Полная аббревиатура — OODA: "
        "Observe — Orient — Decide — Act. Мы её используем как несущую ось "
        "всей лекции, упростив до трёх звеньев: Sense, Decide, Act.\n\n"
        "Главное, что нужно унести с этого слайда: AI входит в каждое звено "
        "по-разному, и провалы случаются преимущественно на стыках звеньев. "
        "В Sense AI работает лучше всего — данных много, ground truth доступна, "
        "цена ошибки терпима. В Decide AI работает как ускоритель аналитика, "
        "но плохо как замена; здесь LLM-галлюцинации и automation bias "
        "масштабируются в крупные ошибки. В Act AI работает в узких сценариях "
        "с supervised pilots overhead, но полная автономия в бою в 2026 году — "
        "это маркетинг и переговорный концепт, не развёрнутая практика.")
    return s


def section_divider(prs, number, title, frame_phrase, current_section, caption):
    """Generic section divider slide."""
    s = blank(prs); set_slide_bg(s, WHITE)

    # Huge decorative number on left
    text_box(s, 0.0, 0.7, 5.5, 6.0, str(number),
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD, line_spacing=0.9)

    # Right: title + frame
    text_box(s, 5.8, 2.5, 7.0, 1.5, title,
             size=36, bold=True, color=DEEP, line_spacing=1.1)
    text_box(s, 5.8, 4.3, 7.0, 1.2, frame_phrase,
             size=20, italic=True, color=MID, line_spacing=1.25)

    # Caption above progress bar
    text_box(s, 0.6, 6.55, 12.13, 0.35, caption,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)

    # Progress bar
    add_progress_bar(s, current_section=current_section, y=6.95)
    return s


def slide_06_section1_divider(prs):
    return section_divider(prs, 1, "Sense — глаза в небе и на земле",
        "Где AI работает лучше всего — и где он ломается даже здесь",
        current_section=1,
        caption="12 минут · 5 рабочих кейсов · 3 провала · 2 критерия")


def slide_07_sense_intro(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Sense — самое благополучное звено OODA: эталон есть, цена FP мала, слияние работает")

    # Left — 4 sensor types в 2×2 grid
    text_box(s, 0.6, 1.6, 7.4, 0.5, "Четыре типа сигнала",
             size=18, bold=True, color=MID)
    sensors = [
        ("satellite", "Спутниковая съёмка", "оптика · мультиспектр · IR · SAR"),
        ("gauge", "Телеметрия платформ", "двигатели · бортовые системы · вибро"),
        ("radio", "Радиосигналы противника", "SIGINT · ELINT · COMINT"),
        ("globe", "Открытые источники", "соцсети · новости · AIS commercial"),
    ]
    sc_w = 3.5; sc_h = 1.85; gap = 0.18
    for i, (icon, title, desc) in enumerate(sensors):
        row = i // 2; col = i % 2
        x = 0.6 + col * (sc_w + gap)
        y = 2.15 + row * (sc_h + gap)
        ocean_box(s, x, y, sc_w, sc_h)
        icon_p = ASSETS / "icons" / f"{icon}-48.png"
        if icon_p.exists():
            add_image(s, icon_p, x + 0.2, y + 0.3, w=0.55, h=0.55)
        text_box(s, x + 0.95, y + 0.3, sc_w - 1.1, 0.5, title,
                 size=14, bold=True, color=DEEP)
        text_box(s, x + 0.3, y + 0.95, sc_w - 0.5, 0.85, desc,
                 size=11, italic=True, color=MID, line_spacing=1.2)

    # Right — 3 reasons stack
    text_box(s, 8.5, 1.6, 4.3, 0.5, "Почему здесь AI работает",
             size=18, bold=True, color=MID)
    reasons = [
        ("1", "Доступная эталонная разметка",
         "Снимки проверяются полевым выездом и историей."),
        ("2", "Терпимая цена ложной тревоги",
         "Лишний сигнал — время аналитика, не жизнь."),
        ("3", "Слияние нескольких сенсоров",
         "Снижает индивидуальную хрупкость моделей.", GOLD),
    ]
    rc_h = 1.4; rc_gap = 0.15
    for i, item in enumerate(reasons):
        if len(item) == 4:
            num, title, desc, accent = item
        else:
            num, title, desc = item; accent = MID
        y = 2.15 + i * (rc_h + rc_gap)
        ocean_box(s, 8.5, y, 4.3, rc_h)
        # Number badge
        text_box(s, 8.65, y + 0.2, 0.6, 0.6, num,
                 size=32, bold=True, color=accent,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 9.3, y + 0.18, 3.4, 0.4, title,
                 size=14, bold=True, color=DEEP)
        text_box(s, 9.3, y + 0.6, 3.4, 0.7, desc,
                 size=11, color=MID, italic=True, line_spacing=1.2)

    add_speaker_notes(s,
        "Sense в нашей цепи — это сбор и первичная интерпретация сигналов о "
        "внешнем мире. В аэрокосмическом и оборонном контексте «сигнал» имеет "
        "четыре основных типа. Первый — спутниковая съёмка: оптика, мультиспектр, "
        "инфракрасное, радар с синтезированной апертурой. Второй — телеметрия с "
        "собственных платформ: данные с двигателей, бортовых систем, расход "
        "топлива, виброметрия. Третий — радиосигналы противника: перехват "
        "сообщений и анализ электромагнитной обстановки. Четвёртый — открытые "
        "источники: социальные сети, новости, коммерческие AIS-данные о судах.\n\n"
        "Почему Sense — самое благополучное звено? Три причины, которые надо "
        "запомнить. Доступная ground truth: снимки спутников можно верифицировать "
        "полевым выездом или вторым сенсором; аномалии двигателей — послеполётным "
        "досмотром. Терпимая FP-цена: лишний alert на спутниковой аналитике — "
        "это потерянное время аналитика, не жизнь. И множество сенсоров: можно "
        "строить fusion — слияние нескольких источников, что снижает "
        "индивидуальную хрупкость каждой модели. Multi-sensor fusion — это "
        "структурное преимущество Sense над Decide и Act, к которому мы будем "
        "возвращаться весь раздел.")
    return s


# ===== Continue with remaining slide builders... =====
# (Будем реализовывать каждый отдельной функцией. Из-за объёма — выноcим
# часть слайдов как обобщённые "card layout" функции.)

def slide_08_maxar_sentry(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Maxar Sentry — упреждающая разведка в часах, не днях", size=26)

    # Left — real satellite imagery
    ocean_box(s, 0.6, 1.8, 7.5, 4.8)
    text_box(s, 0.85, 1.95, 7.0, 0.4, "Спутниковая съёмка + обнаружение изменений",
             size=13, italic=True, color=MID, bold=True)
    photo_p = ASSETS / "photos" / "p01-sat-imagery.jpg"
    if photo_p.exists():
        add_image(s, photo_p, 0.85, 2.5, w=7.0, h=3.6)
        text_box(s, 0.85, 6.15, 7.0, 0.4,
                 "Радарная (SAR) съёмка — пример класса данных Sentry; цифры 1, 2, 3 — отмеченные ИИ изменения. Источник: Wikimedia Commons, CC-BY-SA.",
                 size=10, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right — 4 info cards
    cards_data = [
        ("250 ПБ", "архив", "20+ лет данных", GOLD),
        ("NGA Luno A", "главный контракт", "Детекции от ИИ", MID),
        ("3 сенсора", "слияние", "оптика + радар + AIS судов", LIGHT),
        ("Часы", "после съёмки", "не дни, не недели", GOLD),
    ]
    cc_w = 4.3; cc_h = 1.15; cy = 1.8
    for big, lbl, sub, color in cards_data:
        ocean_box(s, 8.5, cy, cc_w, cc_h, fill=LIGHT_TINT if color != GOLD else GOLD_TINT)
        text_box(s, 8.7, cy + 0.12, 2.4, 0.5, big,
                 size=22, bold=True, color=color)
        text_box(s, 8.7, cy + 0.6, 2.5, 0.3, lbl,
                 size=12, color=DEEP, italic=True)
        text_box(s, 11.2, cy + 0.3, 1.5, 0.6, sub,
                 size=10, color=MID, italic=True, align=PP_ALIGN.RIGHT, line_spacing=1.2)
        cy += cc_h + 0.1

    # Anti-hype caveat — russified
    text_box(s, 0.6, 6.7, 12.13, 0.35,
             "Без преувеличений: «AI-derived» в маркетинге означает оркестрацию классической компьютерной разведки + обнаружение изменений + межсенсорное наведение, а не одну фундаментальную модель.",
             size=11, italic=True, color=LIGHT)

    add_footer(s, "Источники: Defense One июнь 2025 · BusinessWire 20250625")
    add_speaker_notes(s,
        "К 2026 году в коммерческой спутниковой аналитике сложилась устойчивая "
        "четвёрка игроков. Maxar Sentry — самый яркий из них. Sentry запущен "
        "25 июня 2025 года как «predictive intelligence suite». Под капотом — "
        "ML-модели над архивом около 250 петабайт снимков плюс мульти-сенсорное "
        "cross-cueing: электро-оптика, радар с синтезированной апертурой и "
        "автоматическая идентификационная система судов. Главный контракт — "
        "NGA, программа Luno A D01, в рамках которой Maxar обязан выдавать "
        "AI-генерируемые детекции самолётов, кораблей и техники в часах после "
        "съёмки. Anti-hype оговорка важна для инженерного слуха. Бренд "
        "«AI-derived detection» часто означает не один LLM-class model, а "
        "оркестрацию из классических методов компьютерного зрения, change "
        "detection и multi-sensor tipping. Это нормально и инженерно правильно — "
        "маркетинговый язык опережает технический. Когда вы будете оценивать "
        "вакансию или предложение от подобного вендора, помните: за словом «AI» "
        "может стоять и foundation model, и хорошо отлаженный pipeline "
        "классических методов. Это разные инженерные стеки.")
    return s


# (Из-за лимита строк продолжаем во второй части файла. См. build_lec09_part2.py
# который импортируется ниже.)


# ========== Main ==========

def main():
    prs = setup_pres()

    # Phase 6 — building all 43 slides
    slide_01_hook(prs)
    slide_02_cover(prs)
    slide_03_lecture_map(prs)
    slide_04_glossary(prs)
    slide_05_keystone(prs)
    slide_06_section1_divider(prs)
    slide_07_sense_intro(prs)
    slide_08_maxar_sentry(prs)

    # Import rest from part2 — to keep this file under 600 lines
    from build_lec09_part2 import build_part2
    build_part2(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
