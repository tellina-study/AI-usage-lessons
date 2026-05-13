"""
Full 33-slide build of Лекции 1 v3.1 (Phase 12.4 revision of EPIC #64, issue #70).

Source-of-truth: deck.yaml v3.1 + chapter v3.1 (status=reviewed, 16406 слов) +
slides/*.md v3.1 (33 файлов с readable speaker notes 150-300 слов).

v3.1 changes vs v3 (4-critic synthesis 2026-05-13):
- 34 → 33 slides: removed s26 ARC-AGI economics (concept retains poorly при self-study)
  and s28 Pearl 3 levels (концептуально красиво, но к концу 75-min лекции не заходит).
- Renumbered: s27→s26 (4-speaker AGI table), s29→s28 (summary+homework), s30→s29 (roadmap),
  s31→s30 (lec2 teaser), s32→s31 (Q&A).
- Added NEW s27 — section 5 divider («Что забрать домой») per DoD §10 + reader-rendered feedback.
- Critical fixes (4):
  * s13 speaker notes synced with visual (Model = left-top, Agent = right-bottom).
  * «Приложение-робот» renamed to «Приложение (автоматизация)» on s21 quadrant;
    s20+s21 notes explain 2 types of apps (with UI / without UI).
  * s05b funnel «10% в проде» → «10% доходят до прода»; widened gold plate.
  * NEW s27 divider section 5.
- High-value fixes (7):
  * s13 axis labels enlarged (10pt → 13-16pt).
  * s15 RU/EN sub-labels unified to RU.
  * s21 axis labels Q1/Q2 moved INSIDE quadrant.
  * s08 «90% откатов» n=50 caveat added to speaker notes.
  * s07 Vaswani citation timestamp «на май 2026».
  * s29 PARTS disclaimer added to speaker notes.
  * s28 takeaway 3 — removed Pearl reference (Pearl slide deleted).

v3 baseline (preserved): Ocean Gradient palette, Ocean rounded box motif on every content
slide, Gold ≥1×/slide, footer-tax = 0, all notes are readable text 150-300 words.

Canvas: 13.333" × 7.5" (16:9). Pacing: 62.5 active + 12.5 buffer = 75 min.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path("/home/levko/AI-usage-lessons/library/lectures/lec-01")
ASSETS = ROOT / "rendered/assets"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-01.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0, radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.45, h=1.15, w=12.3, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.18, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True):
    box = filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                      radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.08, w=w - 0.4, h=h - 0.16, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


def speaker_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


# ============================================================
# Speaker notes loader from md
# ============================================================
def load_notes(slide_id):
    """Extract Speaker notes + Лектору block from slide markdown."""
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    # Find ## Speaker notes section, then ## Лектору
    notes_match = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)', md, re.DOTALL)
    notes = notes_match.group(1).strip() if notes_match else ""
    lect_match = re.search(r'## Лектору\s*\n(.*?)(?=\n## |\Z)', md, re.DOTALL)
    lect = lect_match.group(1).strip() if lect_match else ""
    full = notes
    if lect:
        full += "\n\n--- ЛЕКТОРУ ---\n" + lect
    return full


# ============================================================
# Roadmap bar (used by section dividers)
# ============================================================
def roadmap_bar(slide, here_idx):
    """Render a 5-section roadmap bar at bottom of slide.
    here_idx: 0=section 0 (open), 1=раздел 1, 2=раздел 2, 3=раздел 3, 4=раздел 4, 5=раздел 5."""
    # 6 cells (0..5) over 12.3 width, with 0.05 gaps
    bar_y = 6.55
    bar_h = 0.4
    n_cells = 6
    total_w = 12.3
    gap = 0.06
    cell_w = (total_w - gap * (n_cells - 1)) / n_cells
    start_x = 0.55
    labels = [
        "0  Открытие",
        "1  AI",
        "2  Сейчас",
        "3  Способы",
        "4  Границы",
        "5  Итог",
    ]
    for i, label in enumerate(labels):
        x = start_x + i * (cell_w + gap)
        is_here = (i == here_idx)
        fill = GOLD if is_here else SOFT_GREY
        text_color = DEEP if is_here else SLATE
        filled_rect(slide, x, bar_y, cell_w, bar_h, fill, radius=True, radius_adj=0.30)
        text_box(slide, x=x, y=bar_y + 0.08, w=cell_w, h=bar_h - 0.16,
                 text=label, size=11, bold=is_here, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, x=0.55, y=bar_y + bar_h + 0.02, w=12.3, h=0.25,
             text=f"Вы здесь — раздел {here_idx} из 5",
             size=10, italic=True, color=GOLD, align=PP_ALIGN.LEFT)


# ============================================================
# Slide builders
# ============================================================
def build_s01(p):
    s = blank(p)
    text_box(s, x=0.55, y=0.55, w=5.9, h=2.4,
             text="Идентификация людей в реальном времени — на ноутбуке, без интернета, с 2023 года.",
             size=26, bold=True, color=DEEP, line_spacing=1.20)
    text_box(s, x=0.55, y=3.15, w=5.9, h=1.4,
             text="Narrow AI — модель решает одну задачу (обнаружение людей в кадре) и больше ничего.",
             size=15, italic=True, color=MID, line_spacing=1.3)
    # Bottom caption with mixed runs
    text_runs(s, 0.55, 5.5, 5.9, 1.0, [
        {"text": "На экране — ", "size": 15, "color": DEEP},
        {"text": "YOLOv8", "size": 15, "color": MID, "bold": True},
        {"text": " на CPU ноутбука: ", "size": 15, "color": DEEP},
        {"text": "~30 fps", "size": 15, "color": GOLD, "bold": True},
        {"text": ".", "size": 15, "color": DEEP},
        {"newpara": True, "text": "Без интернета", "size": 15, "color": TEAL, "bold": True},
        {"text": "  ·  обучена в 2023.", "size": 15, "color": DEEP},
    ], line_spacing=1.35)
    # Right column — Ocean rounded box framing the YOLO mock screenshot
    box_x, box_y, box_w, box_h = 6.55, 0.55, 6.3, 4.4
    ocean_box(s, box_x, box_y, box_w, box_h)
    pad = 0.18
    img_w = box_w - 2 * pad
    img_h = img_w * 720.0 / 1280.0
    img_x = box_x + pad
    img_y = box_y + (box_h - img_h) / 2.0
    add_image(s, ASSETS / "illustrations/s01-yolo-mock.png", img_x, img_y, img_w, img_h)
    text_box(s, x=box_x, y=box_y + box_h + 0.05, w=box_w, h=0.4,
             text="Кадр модели в момент демо: 2 человека в боксах. YOLOv8 (Ultralytics, 2023).",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Cover — distinct: tinted bg, decorative «01», 60pt title."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=8.0, y=2.7, w=5.3, h=4.7, text="01",
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.7, y=1.0, w=6.5, h=0.55, text="ЛЕКЦИЯ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.72, 1.55, 0.7, 0.05, fill=TEAL)
    text_box(s, x=0.7, y=2.0, w=8.5, h=2.4, text="Введение —\nAI вокруг нас",
             size=60, bold=True, color=DEEP, line_spacing=1.05, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.7, 5.45, 0.05, 0.55, fill=TEAL)
    text_box(s, x=0.95, y=5.45, w=8.0, h=0.6,
             text="Карта применений AI: где работает, где — нет.",
             size=22, color=MID, italic=False, align=PP_ALIGN.LEFT, line_spacing=1.25)
    if (ASSETS / "illustrations/hero-cover-light.png").exists():
        add_image(s, ASSETS / "illustrations/hero-cover-light.png",
                  x=8.0, y=0.9, w=5.0, h=5.0)
    speaker_notes(s, load_notes("s02"))


def build_s02a(p):
    """Lecture map — 5 sections horizontal. Fix-12: timing removed."""
    s = blank(p)
    slide_title(s, "Карта лекции — 5 разделов", size=28)
    sections = [
        # (num, title, 1-phrase description, color)
        ("0", "Открытие\nи опросы", "Где мы сейчас\nкак пользователи", LIGHT),
        ("1", "Что такое AI", "Определения,\nистория, перелом", LIGHT),
        ("2", "Где мы\nсейчас", "Цифры рынка\n2022–2026", MID),
        ("3", "Четыре способа\nреализации", "Модель · чат ·\nагент · приложение", MID),
        ("4", "Границы\nи безопасность", "Что AI ломает\nи где не работает", DEEP),
        ("5", "Заключение", "Резюме · задание ·\nкарта семестра", DEEP),
    ]
    card_y = 2.2
    card_w = 1.95
    card_h = 3.0
    gap = 0.15
    start_x = (SLIDE_W_IN - (card_w * 6 + gap * 5)) / 2.0
    for i, (num, title, desc, color) in enumerate(sections):
        x = start_x + i * (card_w + gap)
        is_here = (i == 0)
        fill = SURFACE
        ocean_box(s, x, card_y, card_w, card_h, fill=fill,
                  stroke=GOLD if is_here else LIGHT,
                  stroke_pt=2.5 if is_here else 1.5)
        text_box(s, x=x, y=card_y + 0.25, w=card_w, h=0.6, text=num,
                 size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.1, y=card_y + 1.15, w=card_w - 0.2, h=0.95, text=title,
                 size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.20)
        text_box(s, x=x + 0.1, y=card_y + 2.05, w=card_w - 0.2, h=0.85, text=desc,
                 size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER, line_spacing=1.25)
    # Gold "Вы здесь" pointer (under section 0). Widen box so text fits one line.
    text_box(s, x=start_x - 0.6, y=card_y + card_h + 0.25, w=card_w + 1.2, h=0.4,
             text="↑ Вы здесь — Раздел 0", size=13, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s02a"))


def build_s03(p):
    s = blank(p)
    slide_title(s, "Сначала — ваша оценка, потом — данные.", size=28)
    card_y = 1.95
    card_h = 4.7
    card_w = 5.95
    # Q1
    q1_x = 0.55
    ocean_box(s, q1_x, card_y, card_w, card_h)
    add_image(s, ASSETS / "icons/lucide-hand-blue.png",
              x=q1_x + 0.35, y=card_y + 0.35, w=0.95, h=0.95)
    text_box(s, x=q1_x + 1.5, y=card_y + 0.4, w=card_w - 1.6, h=0.35,
             text="Вопрос 1  ·  один вариант ответа",
             size=14, bold=True, color=MID)
    text_box(s, x=q1_x + 1.5, y=card_y + 0.75, w=card_w - 1.6, h=0.7,
             text="Как часто вы используете AI?",
             size=22, bold=True, color=DEEP, line_spacing=1.20)
    chip_y = card_y + 2.1
    chip_h = 0.5
    chips = ["Никогда", "Несколько раз\nв месяц", "Несколько раз\nв неделю", "Каждый день"]
    chip_w = 1.25
    for i, ctxt in enumerate(chips):
        x = q1_x + 0.35 + i * (chip_w + 0.05)
        chip(s, x, chip_y, chip_w, 0.85, ctxt, fill=MID, color=WHITE, size=11)
    text_box(s, x=q1_x + 0.35, y=card_y + 3.5, w=card_w - 0.7, h=0.85,
             text="(поднимите руку или ответьте в чате — увидим распределение)",
             size=12, italic=True, color=LIGHT, line_spacing=1.30)
    # Q2
    q2_x = q1_x + card_w + 0.35
    ocean_box(s, q2_x, card_y, card_w, card_h)
    add_image(s, ASSETS / "icons/lucide-message-square-blue.png",
              x=q2_x + 0.35, y=card_y + 0.35, w=0.95, h=0.95)
    text_box(s, x=q2_x + 1.5, y=card_y + 0.4, w=card_w - 1.6, h=0.35,
             text="Вопрос 2  ·  можно несколько вариантов",
             size=14, bold=True, color=TEAL)
    text_box(s, x=q2_x + 1.5, y=card_y + 0.75, w=card_w - 1.6, h=0.7,
             text="Какими AI-инструментами\nвы пользовались хотя бы раз?",
             size=18, bold=True, color=DEEP, line_spacing=1.18)
    chips2 = ["ChatGPT", "YandexGPT", "DeepSeek", "GigaChat",
              "Шедеврум", "Claude", "Gemini", "Ничем"]
    cy = card_y + 2.4
    for i, ctxt in enumerate(chips2):
        col = i % 4
        row = i // 4
        x = q2_x + 0.35 + col * 1.32
        y = cy + row * 0.65
        chip(s, x, y, 1.25, 0.5, ctxt, fill=WHITE, color=TEAL,
             stroke=TEAL, size=12)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    s = blank(p)
    slide_title(s, "Разница между вашей оценкой и реальностью — карта ваших слепых зон про AI.", size=24)
    # Left donut — explicit height matching aspect
    donut_x, donut_y, donut_w, donut_h = 0.55, 1.95, 5.6, 4.2
    ocean_box(s, donut_x, donut_y, donut_w, donut_h)
    text_box(s, x=donut_x + 0.3, y=donut_y + 0.25, w=donut_w - 0.6, h=0.5,
             text="51% российских интернет-пользователей 18+ используют AI раз в неделю и чаще",
             size=14, bold=True, color=DEEP, line_spacing=1.25)
    if (ASSETS / "charts/c1-vciom-donut.png").exists():
        # Constrain image to fit inside box: max width 4.0, max height 2.5
        img_h_max = donut_h - 1.6  # ~2.6
        img_w = min(donut_w - 1.6, img_h_max * 1.5)  # donut is roughly square but file may be wide
        img_h = img_w * 0.67  # aspect ratio for the cached chart (600x400)
        if img_h > img_h_max:
            img_h = img_h_max
            img_w = img_h * 1.5
        add_image(s, ASSETS / "charts/c1-vciom-donut.png",
                  x=donut_x + (donut_w - img_w) / 2, y=donut_y + 0.95, w=img_w, h=img_h)
    text_box(s, x=donut_x + 0.3, y=donut_y + donut_h - 0.45, w=donut_w - 0.6, h=0.35,
             text="ВЦИОМ-Онлайн, 13–15 декабря 2025, n=3239",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Right bar chart
    bar_x, bar_y, bar_w, bar_h = donut_x + donut_w + 0.3, 1.95, 6.6, 4.2
    ocean_box(s, bar_x, bar_y, bar_w, bar_h)
    text_box(s, x=bar_x + 0.3, y=bar_y + 0.25, w=bar_w - 0.6, h=0.5,
             text="Использование LLM в РФ среди пользователей AI",
             size=14, bold=True, color=DEEP, line_spacing=1.25)
    # Manual horizontal bars
    rows = [("ChatGPT", 27, MID), ("YandexGPT", 23, MID),
            ("DeepSeek", 20, GOLD), ("GigaChat", 15, MID), ("Шедеврум", 11, MID)]
    bar_top = bar_y + 1.0
    row_h = 0.42
    for i, (label, pct, color) in enumerate(rows):
        ry = bar_top + i * row_h
        text_box(s, x=bar_x + 0.3, y=ry + 0.02, w=1.4, h=0.35, text=label,
                 size=13, bold=True, color=DEEP)
        bar_max_w = 3.6
        bar_actual_w = bar_max_w * pct / 50.0
        filled_rect(s, bar_x + 1.7, ry + 0.08, bar_actual_w, 0.25, color, radius=True, radius_adj=0.5)
        text_box(s, x=bar_x + 1.75 + bar_actual_w, y=ry + 0.05, w=0.7, h=0.35,
                 text=f"{pct}%", size=13, bold=True, color=DEEP if color == GOLD else MID)
    # DeepSeek callout — push below bars with safe gap
    text_box(s, x=bar_x + 0.3, y=bar_top + 5 * row_h + 0.20, w=bar_w - 0.6, h=0.40,
             text="↑ Та же страна, телеметрия Microsoft (2026): DeepSeek = 43%",
             size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    text_box(s, x=bar_x + 0.3, y=bar_y + bar_h - 0.45, w=bar_w - 0.6, h=0.4,
             text="ВЦИОМ окт 2025, n=1600, multi-select. Сумма >100% — респонденты могли указать несколько.",
             size=10, italic=True, color=LIGHT)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.35, 12.25, 0.55,
                 "Сравнивайте методологии прежде, чем сравнивать AI-цифры.")
    speaker_notes(s, load_notes("s04"))


def build_s05a(p):
    s = blank(p)
    slide_title(s, "Кто я и почему мне это важно.", size=28)
    # Left monogram
    mono_x, mono_y, mono_d = 1.0, 2.2, 3.5
    monogram = filled_rect(s, mono_x, mono_y, mono_d, mono_d, MID, radius=True, radius_adj=0.5)
    monogram.line.fill.background()
    text_box(s, x=mono_x, y=mono_y + 0.65, w=mono_d, h=mono_d - 1.0,
             text="ИИ", size=120, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=mono_x - 0.3, y=mono_y + mono_d + 0.1, w=mono_d + 0.6, h=0.4,
             text="инициалы лектора", size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right 3 cards
    cards = [
        ("lucide-briefcase-blue.png", "Опыт", "[годы работы с AI; конкретные проекты]"),
        ("lucide-lightbulb-blue.png", "Зачем мне курс", "[почему важно лично]"),
        ("lucide-users-blue.png", "Связь со студентами", "[контакт, формат вопросов]"),
    ]
    card_x = 5.5
    card_w = 7.4
    card_h = 1.4
    card_gap = 0.15
    card_y_start = 2.0
    for i, (icon, title, body) in enumerate(cards):
        y = card_y_start + i * (card_h + card_gap)
        ocean_box(s, card_x, y, card_w, card_h)
        add_image(s, ASSETS / "icons" / icon, x=card_x + 0.35, y=y + 0.4,
                  w=0.65, h=0.65)
        text_box(s, x=card_x + 1.2, y=y + 0.25, w=card_w - 1.4, h=0.4,
                 text=title, size=16, bold=True, color=MID)
        text_box(s, x=card_x + 1.2, y=y + 0.75, w=card_w - 1.4, h=0.55,
                 text=body, size=13, italic=True, color=SLATE, line_spacing=1.30)
    speaker_notes(s, load_notes("s05a"))


def build_s05b(p):
    s = blank(p)
    slide_title(s, "Главный вопрос курса — не «можно ли AI?», а «нужно ли и где?».", size=24)
    # Left: funnel
    fun_x, fun_y, fun_w = 0.55, 2.0, 5.2
    # Top trapezoid (100%)
    filled_rect(s, fun_x, fun_y, fun_w, 1.05, LIGHT, radius=True, radius_adj=0.08)
    text_box(s, x=fun_x, y=fun_y + 0.18, w=fun_w, h=0.7,
             text="100%  AI-пилотов\nзапускаются",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)
    # Middle (−90%)
    mid_w = fun_w * 0.65
    mid_x = fun_x + (fun_w - mid_w) / 2.0
    filled_rect(s, mid_x, fun_y + 1.20, mid_w, 1.05, MID, radius=True, radius_adj=0.10)
    text_box(s, x=mid_x, y=fun_y + 1.35, w=mid_w, h=0.7,
             text="−90%  откатываются",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom (10%) — gold; widened so «10% доходят до прода» fits on one line
    bot_w = fun_w * 0.85
    bot_x = fun_x + (fun_w - bot_w) / 2.0
    filled_rect(s, bot_x, fun_y + 2.40, bot_w, 1.05, GOLD, radius=True, radius_adj=0.10)
    text_box(s, x=bot_x, y=fun_y + 2.55, w=bot_w, h=0.7,
             text="10% доходят до прода",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=fun_x, y=fun_y + 3.6, w=fun_w, h=0.5,
             text="Иллюстрация принципа (Gartner, McKinsey подтверждают похожие цифры).",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Right: takeaway + central question
    right_x = 6.3
    right_w = 6.5
    ocean_box(s, right_x, fun_y - 0.3, right_w, 5.0)
    text_box(s, x=right_x + 0.3, y=fun_y, w=right_w - 0.6, h=0.45,
             text="Главная мысль лекции",
             size=14, bold=True, color=TEAL)
    text_box(s, x=right_x + 0.3, y=fun_y + 0.5, w=right_w - 0.6, h=1.6,
             text="Завтра — почти везде.\nСегодня — почти никто.\nКурс — про этот разрыв.",
             size=22, bold=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=right_x + 0.3, y=fun_y + 2.4, w=right_w - 0.6, h=0.45,
             text="Центральный вопрос курса",
             size=14, bold=True, color=GOLD)
    text_box(s, x=right_x + 0.3, y=fun_y + 2.85, w=right_w - 0.6, h=1.6,
             text="Где AI работает,\nгде — нет,\nи как это понять?",
             size=24, bold=True, color=DEEP, line_spacing=1.25)
    speaker_notes(s, load_notes("s05b"))


def build_s06(p):
    """Multiple definitions of AI — 4 approaches grid (full definitions, not labels) + AI Effect.

    Fix-13 (2026-05-13): cards now contain the actual definitions (~15-20 words each)
    instead of just approach names. Body 14pt (was 12pt), source 11pt italic (was 10pt),
    cell_h 2.40" (was 1.95") to fit full text. Grid moved up; callout moved down.
    """
    s = blank(p)
    slide_title(s, "Определений AI много — потому что AI это moving target.", size=26)
    cards = [
        ("Russell & Norvig (AIMA, 2021)",
         "«AI = система, мыслящая как человек, мыслящая рационально, действующая как человек или действующая рационально (4 квадранта по 2 осям).»",
         "Russell & Norvig, AIMA, 4-е изд., 2021",
         MID),
        ("ISO/IEC 22989:2022",
         "«AI-система — это engineered system, которая генерирует выходы (рекомендации, прогнозы, решения) для целей, заданных человеком.»",
         "Международный стандарт ISO/IEC 22989:2022 — опора EU AI Act",
         LIGHT),
        ("Через обучение (Mitchell, 1997)",
         "«Программа улучшается с опытом E на задаче T по метрике P. Если поведение возникает из обученной модели — это AI.»",
         "Mitchell, Machine Learning, 1997",
         MID),
        ("Через бенчмарки и AGI",
         "«AI = то, что проходит тест Тьюринга или решает бенчмарк на уровне человека.» Возражение Сёрла: поведение ≠ понимание.",
         "Тьюринг 1950 / Searle 1980 — Chinese Room",
         LIGHT),
    ]
    grid_x = 0.55
    grid_y = 1.62
    cell_w = 6.05
    cell_h = 2.40
    cell_gap = 0.15
    for i, (head, body, src, color) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = grid_x + col * (cell_w + cell_gap)
        y = grid_y + row * (cell_h + cell_gap)
        ocean_box(s, x, y, cell_w, cell_h, stroke=color)
        # Header (approach name)
        text_box(s, x=x + 0.25, y=y + 0.18, w=cell_w - 0.5, h=0.40,
                 text=head, size=15, bold=True, color=color)
        # Body — full definition, 14pt for projector readability
        text_box(s, x=x + 0.25, y=y + 0.60, w=cell_w - 0.5, h=1.40,
                 text=body, size=14, color=DEEP, line_spacing=1.22)
        # Source citation
        text_box(s, x=x + 0.25, y=y + cell_h - 0.40, w=cell_w - 0.5, h=0.32,
                 text=src, size=11, italic=True, color=SLATE)
    # AI Effect callout at bottom (gold accent, ≥1×/slide rule)
    gold_callout(s, 0.55, 6.85, 12.25, 0.55,
                 "AI Effect (Tesler):  «AI is whatever hasn't been done yet».  Как только техника начинает работать, её перестают называть AI.",
                 size=13)
    speaker_notes(s, load_notes("s06"))


def build_s07(p):
    """70 years AI timeline — 3 stacked group rows, no overlap."""
    s = blank(p)
    slide_title(s, "70 лет AI: открытия, зимы, точка перелома 2017.", size=28)
    # Three group bands stacked vertically — each row contains its own events
    groups = [
        ("Открытия (1950 — 1980-е)", LIGHT, [
            ("1950", "Turing\nImitation Game"),
            ("1956", "Дартмут\nMcCarthy вводит «AI»"),
            ("1966", "ELIZA\nWeizenbaum"),
            ("1980-е", "Экспертные\nсистемы"),
        ]),
        ("Зимы и прорывы (1973 — 2012)", MID, [
            ("1974", "1-я зима\nLighthill"),
            ("1987", "2-я зима\nкрах Lisp Machines"),
            ("1997", "Deep Blue\n200M поз/сек"),
            ("2012", "AlexNet\nGPU + DL"),
        ]),
        ("Перелом и взрыв (2017 — 2026)", DEEP, [
            ("2017", "«Attention Is\nAll You Need»  ★"),
            ("2022", "ChatGPT\n1M за 5 дней"),
            ("2024", "MCP\nде-факто стандарт"),
            ("2025-26", "DeepSeek R1,\nClaude Code"),
        ]),
    ]
    band_h = 1.30
    band_y_start = 1.85
    for gi, (gname, color, events) in enumerate(groups):
        band_y = band_y_start + gi * (band_h + 0.10)
        # Group label (left)
        text_box(s, x=0.55, y=band_y + 0.20, w=2.5, h=0.4, text=gname,
                 size=11, bold=True, color=color)
        # Group color band
        filled_rect(s, 3.20, band_y + 0.55, 9.55, 0.15, color, radius=True, radius_adj=0.5)
        # Event boxes spaced over the band
        n = len(events)
        ev_w = 9.55 / n - 0.15
        for ei, (year, label) in enumerate(events):
            ex = 3.20 + ei * (ev_w + 0.15) + 0.05
            ey = band_y + 0.10
            is_pivot = "★" in label
            # Year tick
            tick_x = ex + ev_w / 2 - 0.10
            tick_y = band_y + 0.55
            if is_pivot:
                shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(tick_x - 0.05), Inches(tick_y - 0.05),
                                         Inches(0.30), Inches(0.30))
                shp.fill.solid(); shp.fill.fore_color.rgb = GOLD
                shp.line.color.rgb = DEEP; shp.line.width = Pt(1.5)
                disable_shadow(shp)
            else:
                filled_rect(s, tick_x + 0.05, tick_y - 0.05, 0.10, 0.25, color)
            # Year label
            text_box(s, x=ex, y=band_y + 0.85, w=ev_w, h=0.30, text=year,
                     size=11, bold=True,
                     color=GOLD if is_pivot else color, align=PP_ALIGN.CENTER)
            # Event label above year tick
            text_box(s, x=ex, y=band_y - 0.05, w=ev_w, h=0.55, text=label,
                     size=9.5, color=DEEP if not is_pivot else GOLD,
                     bold=is_pivot, align=PP_ALIGN.CENTER, line_spacing=1.20)
    # AI Effect callout at bottom — give it room
    gold_callout(s, 0.55, 6.05, 12.25, 0.95,
                 "AI Effect (Tesler):  «AI is whatever hasn't been done yet». Распознавание речи, навигация по пробкам, разблокировка лица — всё это «было AI», теперь — «функция приложения».",
                 size=12)
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """Scale numbers — 4 metrics grid + counter-fact gold."""
    s = blank(p)
    slide_title(s, "AI стал инфраструктурой за 3 года: 900M пользователей, 51% разработчиков ежедневно, 46% Copilot-кода.", size=22)
    metrics = [
        ("900M", "WAU", "ChatGPT, февраль 2026", "OpenAI", MID, "lucide-users-2-blue.png"),
        ("51%", "ежедневно", "Stack Overflow Dev Survey 2025", "n=49k+, 177 стран", LIGHT, "lucide-code-blue.png"),
        ("46%", "кода у Copilot", "GitHub Octoverse 2025", "Java — 61%", MID, "lucide-github.png"),
        ("$244–390B", "AI-рынок", "Statista / McKinsey 2025", "разброс по методологии", LIGHT, "lucide-dollar-sign-blue.png"),
    ]
    grid_y = 2.0
    cell_w = 6.05
    cell_h = 1.85
    grid_x = 0.55
    cell_gap = 0.15
    for i, (big, label, src1, src2, color, icon) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = grid_x + col * (cell_w + cell_gap)
        y = grid_y + row * (cell_h + cell_gap)
        ocean_box(s, x, y, cell_w, cell_h, stroke=color)
        # Icon top right
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + cell_w - 0.85, y=y + 0.25, w=0.55, h=0.55)
        # Big number
        text_box(s, x=x + 0.30, y=y + 0.20, w=cell_w - 1.4, h=0.85,
                 text=big, size=44, bold=True, color=color, line_spacing=1.0)
        text_box(s, x=x + 0.30, y=y + 0.95, w=cell_w - 0.5, h=0.4,
                 text=label, size=15, bold=True, color=DEEP)
        text_box(s, x=x + 0.30, y=y + cell_h - 0.55, w=cell_w - 0.5, h=0.32,
                 text=src1, size=10, italic=True, color=SLATE)
        text_box(s, x=x + 0.30, y=y + cell_h - 0.28, w=cell_w - 0.5, h=0.28,
                 text=src2, size=10, italic=True, color=SLATE)
    # Counter-fact gold strip
    gold_callout(s, 0.55, 6.05, 12.25, 0.90,
                 "Контр-факт: ~90% AI-пилотов в РФ не доходят до прода. CNews / Vedomosti / Intellectual Analytics, март 2026.",
                 size=14)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """4 breakthroughs 2023-2026 — horizontal episodes."""
    s = blank(p)
    slide_title(s, "Пространство открыто: 4 прорыва 2023–2026 от не-первых игроков.", size=26)
    episodes = [
        ("сентябрь\n2023", "Mistral 7B", "Apache 2.0\nобходит Llama-2 13B", "Mistral AI (FR)", MID, False),
        ("апрель\n2024", "Llama-3", "MMLU 79.5\nоткрытые веса", "Meta", LIGHT, False),
        ("январь\n2025", "DeepSeek R1", "$589B\nNvidia drop за день", "DeepSeek (CN)", GOLD, True),
        ("ноябрь\n2024", "MCP", "де-факто стандарт\nAI-интеграций", "Anthropic", MID, False),
    ]
    card_y = 2.05
    card_w = 2.95
    card_h = 4.0
    gap = 0.25
    start_x = (SLIDE_W_IN - (card_w * 4 + gap * 3)) / 2.0
    for i, (date, name, fact, org, color, is_gold) in enumerate(episodes):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h,
                  stroke=color, stroke_pt=2.5 if is_gold else 1.5)
        # Date band top
        filled_rect(s, x, card_y, card_w, 0.7, color, radius=True, radius_adj=0.1)
        text_box(s, x=x, y=card_y + 0.06, w=card_w, h=0.6, text=date,
                 size=14, bold=True, color=DEEP if color == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        # Name big
        text_box(s, x=x + 0.15, y=card_y + 0.95, w=card_w - 0.3, h=0.7, text=name,
                 size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.10)
        # Main fact
        text_box(s, x=x + 0.20, y=card_y + 1.85, w=card_w - 0.4, h=1.4, text=fact,
                 size=14, bold=is_gold, color=color if is_gold else DEEP,
                 align=PP_ALIGN.CENTER, line_spacing=1.30)
        # Org
        text_box(s, x=x, y=card_y + card_h - 0.55, w=card_w, h=0.4, text=org,
                 size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.50, 12.25, 0.55,
                 "Не отчаивайтесь: серьёзные прорывы делают разные команды. Курс — про устойчивые концепты, переживающие смену поколений моделей.",
                 size=13)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    """Section 3 divider — split «Раздел 3» so it doesn't overflow vertically."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # «Раздел» on its own line, smaller, top
    text_box(s, x=0.55, y=0.5, w=12.25, h=1.5, text="Раздел  3",
             size=110, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Title at safe vertical band
    text_box(s, x=0.55, y=3.0, w=12.25, h=1.4,
             text="Четыре способа реализации систем с AI",
             size=40, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    text_box(s, x=0.55, y=4.50, w=12.25, h=0.6,
             text="Не альтернативы, а слои.",
             size=22, italic=True, color=MID, align=PP_ALIGN.CENTER)
    roadmap_bar(s, 3)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """Layers not alternatives — 4 nested rounded boxes (label OUTSIDE for clarity)."""
    s = blank(p)
    slide_title(s, "Способы реализации систем с AI: не альтернативы, а слои.", size=26)
    # Nested layers (right side); labels go OUTSIDE the boxes (top-right corner each)
    cx, cy = 8.7, 4.05
    sizes = [
        # (w, h, color, label) — outermost first
        (5.4, 4.4, DEEP, "Приложение"),
        (4.3, 3.5, MID, "Агент"),
        (3.2, 2.6, LIGHT, "Чат"),
        (2.1, 1.7, TEAL, "Модель"),
    ]
    # Draw boxes first (outer→inner)
    for i, (w, h, color, label) in enumerate(sizes):
        x = cx - w / 2
        y = cy - h / 2
        ocean_box(s, x, y, w, h, fill=WHITE, stroke=color, stroke_pt=2.5)
    # Then labels — top-right area of each box, outside the next inner one
    label_offsets = [
        # Where label sits (relative to box top): label x relative to right edge
        (sizes[0], 0.20),  # outer
        (sizes[1], 0.20),
        (sizes[2], 0.20),
        (sizes[3], 0.20),  # innermost — center
    ]
    for i, (size_tuple, _) in enumerate(label_offsets):
        w, h, color, label = size_tuple
        x = cx - w / 2
        y = cy - h / 2
        if i < 3:
            # Label at top-right of band (inside outer ring of this box)
            text_box(s, x=x + 0.20, y=y + 0.12, w=w - 0.4, h=0.4, text=label,
                     size=13, bold=True, color=color, align=PP_ALIGN.LEFT)
        else:
            # Innermost — center label
            text_box(s, x=x, y=y + h/2 - 0.20, w=w, h=0.4, text=label,
                     size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Left side — explanation
    text_box(s, x=0.55, y=2.0, w=4.5, h=0.55,
             text="Каждый следующий слой",
             size=18, bold=True, color=DEEP)
    text_box(s, x=0.55, y=2.55, w=4.5, h=0.55,
             text="включает предыдущий",
             size=18, bold=True, color=GOLD)
    text_box(s, x=0.55, y=3.30, w=4.5, h=2.5,
             text="Модель — базовый компонент.\nЧат — модель в обёртке.\nАгент — чат + инструменты.\nПриложение — продукт, в котором AI один из компонентов.",
             size=13, color=DEEP, line_spacing=1.50)
    gold_callout(s, 0.55, 6.10, 4.5, 0.85,
                 "Выбор слоя — инженерное решение, не альтернатива.",
                 size=12)
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """Classification matrix — task × modality 2D grid."""
    s = blank(p)
    slide_title(s, "Классификация AI-систем — две оси: тип задачи × модальность.", size=26)
    # Big matrix area
    matrix_x, matrix_y = 0.85, 1.95
    matrix_w, matrix_h = 11.6, 4.6
    ocean_box(s, matrix_x, matrix_y, matrix_w, matrix_h)
    # Tasks (X axis)
    tasks = ["Класси-\nфикация", "Распо-\nзнавание", "Поиск", "Гене-\nрация", "Прогноз", "Плани-\nрование"]
    # Modalities (Y axis)
    modalities = ["Текст", "Изображение", "Звук / видео", "Структ. данные", "Код"]
    # Grid layout
    grid_left = matrix_x + 1.6
    grid_top = matrix_y + 0.85
    grid_w = matrix_w - 1.8
    grid_h = matrix_h - 1.1
    cell_w = grid_w / len(tasks)
    cell_h = grid_h / len(modalities)
    # Header row (tasks)
    for i, t in enumerate(tasks):
        x = grid_left + i * cell_w
        text_box(s, x=x, y=matrix_y + 0.25, w=cell_w, h=0.55, text=t,
                 size=11, bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.10)
    # Header column (modalities)
    for j, m in enumerate(modalities):
        y = grid_top + j * cell_h
        text_box(s, x=matrix_x + 0.20, y=y + cell_h / 2 - 0.15, w=1.4, h=0.4, text=m,
                 size=12, bold=True, color=MID, align=PP_ALIGN.RIGHT)
    # Grid lines
    for i in range(len(tasks) + 1):
        x = grid_left + i * cell_w
        filled_rect(s, x - 0.005, grid_top, 0.01, grid_h, SOFT_GREY)
    for j in range(len(modalities) + 1):
        y = grid_top + j * cell_h
        filled_rect(s, grid_left, y - 0.005, grid_w, 0.01, SOFT_GREY)
    # Place 3 examples in cells
    examples = [
        # (task_idx, modality_idx, label, color)
        (3, 0, "Google\nTranslate", TEAL),    # Generation × Text
        (4, 3, "AlphaFold", MID),              # Forecast × Structured
        (1, 1, "YOLO ←", GOLD),                # Recognition × Image (with gold callback)
    ]
    for ti, mi, label, color in examples:
        x = grid_left + ti * cell_w + 0.1
        y = grid_top + mi * cell_h + 0.15
        cw = cell_w - 0.2
        ch = cell_h - 0.30
        filled_rect(s, x, y, cw, ch, color, radius=True, radius_adj=0.18)
        text_box(s, x=x, y=y + ch/2 - 0.25, w=cw, h=0.55, text=label,
                 size=12, bold=True, color=DEEP if color == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    # Caption: Yolo callback
    text_box(s, x=matrix_x + 0.30, y=matrix_y + matrix_h - 0.55, w=matrix_w - 0.6, h=0.35,
             text="YOLO в gold — был в демо в начале лекции. Translate, AlphaFold — другие примеры из своих ячеек.",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Bottom note
    gold_callout(s, 0.55, 6.65, 12.25, 0.45,
                 "Подход к обучению и архитектура — не основная цель этой лекции; глубже на лекции 2.",
                 size=12)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """Control quadrant 2x2 — model/chat/agent placement."""
    s = blank(p)
    slide_title(s, "Одна задача, три способа: контроль распределяется между разработчиком и пользователем.", size=22)
    # Quadrant area — shrunk vertically a bit to leave room for axis label + callout below
    qx, qy = 1.7, 1.9
    qw, qh = 7.4, 3.95
    # Box outline
    filled_rect(s, qx, qy, qw, qh, WHITE, stroke=LIGHT, stroke_pt=1.5, radius=True, radius_adj=0.04)
    # Internal cross lines
    filled_rect(s, qx, qy + qh / 2 - 0.005, qw, 0.01, SOFT_GREY)
    filled_rect(s, qx + qw / 2 - 0.005, qy, 0.01, qh, SOFT_GREY)
    # Y axis label (left, vertical conceptual) — enlarged per Fix-6, wider box to avoid wrap
    text_box(s, x=qx - 1.65, y=qy + qh / 2 - 0.40, w=1.55, h=0.85,
             text="Контроль\nпользователя",
             size=15, bold=True, color=MID, align=PP_ALIGN.RIGHT, line_spacing=1.18)
    # ↑ arrow + «высокий» at top of Y axis (just outside quadrant, near top-left corner)
    text_box(s, x=qx - 1.45, y=qy - 0.08, w=1.35, h=0.30, text="высокий ↑",
             size=11, bold=True, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # «низкий» — at bottom-left of Y axis area, OUTSIDE quadrant
    text_box(s, x=qx - 1.45, y=qy + qh + 0.08, w=1.35, h=0.28, text="низкий ↓",
             size=11, bold=True, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # X axis label (bottom) — enlarged per Fix-6; «низкий»/«высокий» moved UP to top inside (away from agent)
    text_box(s, x=qx + qw / 2 - 2.0, y=qy + qh + 0.08, w=4.0, h=0.40,
             text="Контроль разработчика",
             size=15, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # Arrow + range markers on TOP edge inside the quadrant (away from circle sub-labels at bottom)
    text_box(s, x=qx + 0.20, y=qy + 0.05, w=1.2, h=0.28, text="← низкий",
             size=10, bold=True, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    text_box(s, x=qx + qw - 1.40, y=qy + 0.05, w=1.2, h=0.28, text="высокий →",
             size=10, bold=True, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # Three points on quadrant
    pts = [
        # (x_in_quad_frac, y_in_quad_frac, label, sub, color, is_gold)
        (0.20, 0.25, "Модель", "Низкий разраб,\nвысокий user", LIGHT, False),
        (0.50, 0.50, "Чат", "Средний / средний", MID, False),
        (0.80, 0.78, "Агент", "Высокий разраб,\nнизкий user", GOLD, True),
    ]
    for fx, fy, label, sub, color, is_gold in pts:
        cx = qx + fx * qw
        cy = qy + fy * qh
        d = 1.05
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - d/2), Inches(cy - d/2),
                                 Inches(d), Inches(d))
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.color.rgb = DEEP; shp.line.width = Pt(2.0 if is_gold else 1.0)
        disable_shadow(shp)
        text_box(s, x=cx - 0.6, y=cy - 0.20, w=1.2, h=0.4, text=label,
                 size=14, bold=True, color=DEEP if is_gold else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=cx - 0.85, y=cy + 0.55, w=1.7, h=0.6, text=sub,
                 size=10, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.20)
    # Right side — fixed task box
    task_x = qx + qw + 0.5
    task_w = SLIDE_W_IN - task_x - 0.35
    ocean_box(s, task_x, qy + 0.5, task_w, 3.3)
    text_box(s, x=task_x + 0.25, y=qy + 0.65, w=task_w - 0.5, h=0.4,
             text="Одна и та же задача", size=14, bold=True, color=TEAL)
    text_box(s, x=task_x + 0.25, y=qy + 1.10, w=task_w - 0.5, h=2.0,
             text="Извлечь поля из входящего PDF-договора:\n• дата подписания\n• контрагент\n• сумма\n• срок действия\n\nи положить в таблицу.",
             size=13, color=DEEP, line_spacing=1.40)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.55, 12.25, 0.55,
                 "Распределение контроля — инженерное решение, а не достоинство одного способа над другим.",
                 size=13)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """Mini-divider: 'Разберём подробнее'."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=2.0, w=12.25, h=1.2,
             text="Разберём подробнее каждый из четырёх типов",
             size=44, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    # 4 icon cards horizontal
    icons = [
        ("lucide-cpu-blue.png", "Модель", "stateless inference", LIGHT, True),
        ("lucide-message-circle-blue.png", "Чат", "+ интерфейс + память", MID, False),
        ("lucide-bot-blue.png", "Агент", "+ инструменты + план", MID, False),
        ("lucide-layout-grid-blue.png", "Приложение", "AI как функция продукта", DEEP, False),
    ]
    card_y = 4.0
    card_w = 2.8
    card_h = 2.0
    gap = 0.30
    start_x = (SLIDE_W_IN - (card_w * 4 + gap * 3)) / 2.0
    for i, (icon, name, sub, color, is_now) in enumerate(icons):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h,
                  stroke=GOLD if is_now else color,
                  stroke_pt=2.5 if is_now else 1.5)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + card_w/2 - 0.4, y=card_y + 0.20, w=0.8, h=0.8)
        text_box(s, x=x, y=card_y + 1.1, w=card_w, h=0.5, text=name,
                 size=18, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        text_box(s, x=x, y=card_y + 1.55, w=card_w, h=0.4, text=sub,
                 size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
        if is_now:
            text_box(s, x=x, y=card_y - 0.40, w=card_w, h=0.4, text="↓ Сейчас сюда",
                     size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    """Model with pipeline schema — 5 horizontal blocks + 4 examples."""
    s = blank(p)
    slide_title(s, "Модель — компонент, не система. Inference: вход → препроцессинг → модель → постпроцессинг → выход.", size=20)
    # Horizontal pipeline
    pip_y = 2.2
    pip_h = 1.6
    blocks = [
        ("Сырой\nвход", "кадр камеры,\nтекст, звук", LIGHT),
        ("Препро-\nцессинг", "масштабирование,\nобрезка, токенизация", LIGHT),
        ("Модель", "инференс", MID),
        ("Постпро-\nцессинг", "фильтрация,\nнормализация", LIGHT),
        ("Выход", "JSON, метка,\nдействие", LIGHT),
    ]
    n = len(blocks)
    block_w = 2.0
    arrow_w = 0.45
    total_w = block_w * n + arrow_w * (n - 1)
    start_x = (SLIDE_W_IN - total_w) / 2.0
    for i, (name, sub, color) in enumerate(blocks):
        x = start_x + i * (block_w + arrow_w)
        is_model = (i == 2)
        filled_rect(s, x, pip_y, block_w, pip_h, color,
                    stroke=DEEP if is_model else None, stroke_pt=2.0 if is_model else 0.0,
                    radius=True, radius_adj=0.15)
        text_box(s, x=x, y=pip_y + 0.18, w=block_w, h=0.7, text=name,
                 size=16 if is_model else 14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.10)
        text_box(s, x=x + 0.10, y=pip_y + 0.95, w=block_w - 0.2, h=0.6, text=sub,
                 size=10, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.20)
        if i < n - 1:
            ax = x + block_w + 0.05
            ay = pip_y + pip_h / 2 - 0.18
            filled_rect(s, ax, ay + 0.13, arrow_w - 0.10, 0.10, GOLD)
            # Arrow head
            shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                     Inches(ax + arrow_w - 0.18), Inches(ay - 0.05),
                                     Inches(0.25), Inches(0.45))
            shp.rotation = 90
            shp.fill.solid(); shp.fill.fore_color.rgb = GOLD
            shp.line.fill.background()
            disable_shadow(shp)
    # Owner labels
    text_box(s, x=start_x, y=pip_y + pip_h + 0.10, w=block_w, h=0.4,
             text="↑ внешняя система", size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    text_box(s, x=start_x + (block_w + arrow_w), y=pip_y + pip_h + 0.10, w=block_w, h=0.4,
             text="↑ разработчик", size=10, italic=True, color=MID, align=PP_ALIGN.CENTER, bold=True)
    text_box(s, x=start_x + 2 * (block_w + arrow_w), y=pip_y + pip_h + 0.10, w=block_w, h=0.4,
             text="↑ AI-модель", size=10, italic=True, color=DEEP, align=PP_ALIGN.CENTER, bold=True)
    text_box(s, x=start_x + 3 * (block_w + arrow_w), y=pip_y + pip_h + 0.10, w=block_w, h=0.4,
             text="↑ разработчик", size=10, italic=True, color=MID, align=PP_ALIGN.CENTER, bold=True)
    text_box(s, x=start_x + 4 * (block_w + arrow_w), y=pip_y + pip_h + 0.10, w=block_w, h=0.4,
             text="↑ приложение", size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # 4 model examples
    examples = [
        ("YOLOv8", "детекция\nна изображениях"),
        ("Whisper", "распознавание\nречи"),
        ("Stable Diffusion", "генерация\nизображений"),
        ("AlphaFold", "прогноз\nструктур белков"),
    ]
    ex_y = 4.95
    ex_w = 2.8
    ex_h = 1.3
    ex_gap = 0.20
    ex_start_x = (SLIDE_W_IN - (ex_w * 4 + ex_gap * 3)) / 2.0
    for i, (name, role) in enumerate(examples):
        x = ex_start_x + i * (ex_w + ex_gap)
        ocean_box(s, x, ex_y, ex_w, ex_h)
        text_box(s, x=x, y=ex_y + 0.20, w=ex_w, h=0.45, text=name,
                 size=15, bold=True, color=MID, align=PP_ALIGN.CENTER)
        text_box(s, x=x, y=ex_y + 0.65, w=ex_w, h=0.6, text=role,
                 size=11, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.25)
    gold_callout(s, 0.55, 6.50, 12.25, 0.55,
                 "Препроцессинг и постпроцессинг — ответственность разработчика. YOLO = 50 строк; рабочая система с YOLO = сотни строк.",
                 size=12)
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    """Chat cycle — 6-step loop. Cycle in center column; callouts in side columns."""
    s = blank(p)
    slide_title(s, "Как работает чат: 6 шагов цикла.", size=28)
    # Center loop — smaller radius, narrower band; pushed right of left callout
    cx, cy = 6.65, 4.30
    radius = 1.85
    import math
    steps = [
        "1. Ввод\nпользователя",
        "2. Сборка\nконтекста",
        "3. Передача\nв LLM",
        "4. Генерация\nответа",
        "5. Добавление\nк истории",
        "6. Показ\nпользователю",
    ]
    n = len(steps)
    for i, step in enumerate(steps):
        angle = -math.pi / 2 + i * 2 * math.pi / n
        bw, bh = 1.55, 0.85
        bx = cx + radius * math.cos(angle) - bw/2
        by = cy + radius * math.sin(angle) - bh/2
        ocean_box(s, bx, by, bw, bh, fill=WHITE, stroke=MID)
        text_box(s, x=bx, y=by + 0.10, w=bw, h=bh - 0.20, text=step,
                 size=10, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Center loop indicator — gold circle
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(cx - 0.45), Inches(cy - 0.45),
                             Inches(0.9), Inches(0.9))
    shp.fill.solid(); shp.fill.fore_color.rgb = GOLD
    shp.line.color.rgb = DEEP; shp.line.width = Pt(1.5)
    disable_shadow(shp)
    text_box(s, x=cx - 0.5, y=cy - 0.30, w=1.0, h=0.6, text="LOOP",
             size=12, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Two callouts — placed in upper-left and upper-right corners (above cycle, not overlapping)
    gold_callout(s, 0.55, 1.95, 3.4, 1.5,
                 "Контроль через системный промпт — инженерный рычаг разработчика.",
                 size=12)
    gold_callout(s, 9.40, 1.95, 3.4, 1.5,
                 "Ограничение — контекстное окно: 128k–1M токенов. Старое выпадает.",
                 size=12)
    # Bottom takeaway
    text_box(s, x=0.55, y=6.65, w=12.25, h=0.5,
             text="Чат не «помнит» ничего вне контекстного окна. Чат — это конвейер «собрать → подать → дописать → показать».",
             size=13, italic=True, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """Chat = model + UI + memory; case + LLM bar chart."""
    s = blank(p)
    slide_title(s, "Чат = модель + интерфейс + память диалога.", size=28)
    # Left: case card
    case_x, case_y, case_w, case_h = 0.55, 1.95, 6.5, 4.5
    ocean_box(s, case_x, case_y, case_w, case_h)
    text_box(s, x=case_x + 0.30, y=case_y + 0.25, w=case_w - 0.6, h=0.4,
             text="Кейс — типовой для чата", size=14, bold=True, color=TEAL)
    text_box(s, x=case_x + 0.30, y=case_y + 0.7, w=case_w - 0.6, h=0.95,
             text="Инженер получил непонятное ТЗ от смежного отдела, нужно составить чек-лист своей работы.",
             size=15, color=DEEP, line_spacing=1.30)
    # Mock dialog
    dlg_y = case_y + 1.85
    dlg_h = 2.0
    filled_rect(s, case_x + 0.4, dlg_y, case_w - 0.8, dlg_h, WHITE,
                stroke=SOFT_GREY, stroke_pt=1.0, radius=True, radius_adj=0.08)
    text_runs(s, case_x + 0.6, dlg_y + 0.15, case_w - 1.2, dlg_h - 0.3, [
        {"text": "Вы:  ", "size": 12, "bold": True, "color": MID},
        {"text": "Объясни пункт 4.2 ТЗ простыми словами.", "size": 12, "color": DEEP},
        {"newpara": True, "text": "Чат:  ", "size": 12, "bold": True, "color": TEAL},
        {"text": "Это требование к…", "size": 12, "color": DEEP},
        {"newpara": True, "text": "Вы:  ", "size": 12, "bold": True, "color": MID},
        {"text": "Составь чек-лист на 5 пунктов.", "size": 12, "color": DEEP},
        {"newpara": True, "text": "Чат:  ", "size": 12, "bold": True, "color": TEAL},
        {"text": "1. Проверить…  2. Согласовать…  3. …", "size": 12, "color": DEEP},
    ], line_spacing=1.45)
    text_box(s, x=case_x + 0.30, y=case_y + case_h - 0.5, w=case_w - 0.6, h=0.35,
             text="Разовая задача с уточнениями — оптимально для чата. Не модель, не агент, не приложение.",
             size=11, italic=True, color=LIGHT, line_spacing=1.30)
    # Right: bar chart of LLM shares
    bar_x, bar_y, bar_w, bar_h = case_x + case_w + 0.35, 1.95, 5.4, 4.5
    ocean_box(s, bar_x, bar_y, bar_w, bar_h)
    text_box(s, x=bar_x + 0.25, y=bar_y + 0.25, w=bar_w - 0.5, h=0.4,
             text="Какие LLM используют в РФ", size=14, bold=True, color=DEEP)
    text_box(s, x=bar_x + 0.25, y=bar_y + 0.70, w=bar_w - 0.5, h=0.4,
             text="ВЦИОМ окт 2025, n=1600, multi-select", size=10, italic=True, color=LIGHT)
    rows = [("ChatGPT", 27, MID), ("YandexGPT", 23, MID),
            ("DeepSeek", 20, TEAL), ("GigaChat", 15, MID), ("Шедеврум", 11, MID)]
    bt = bar_y + 1.30
    rh = 0.55
    for i, (label, pct, color) in enumerate(rows):
        ry = bt + i * rh
        text_box(s, x=bar_x + 0.25, y=ry + 0.05, w=1.4, h=0.4, text=label,
                 size=13, bold=True, color=DEEP)
        bar_max = 3.2
        bar_actual = bar_max * pct / 30.0
        filled_rect(s, bar_x + 1.55, ry + 0.10, bar_actual, 0.30, color, radius=True, radius_adj=0.5)
        text_box(s, x=bar_x + 1.6 + bar_actual, y=ry + 0.05, w=0.7, h=0.4,
                 text=f"{pct}%", size=12, bold=True, color=color)
    gold_callout(s, 0.55, 6.55, 12.25, 0.50,
                 "Возвращаемся к: где AI работает, а где — нет?",
                 size=13)
    speaker_notes(s, load_notes("s17"))


def build_s18(p):
    """Agent architecture — Chat + Orchestrator + Memory + Tools."""
    s = blank(p)
    slide_title(s, "Агент = чат + оркестратор + внешняя память + инструменты.", size=26)
    # Center: LLM/Chat
    cx, cy = 6.65, 4.0
    chat_d = 2.4
    chat_h = 1.6
    filled_rect(s, cx - chat_d/2, cy - chat_h/2, chat_d, chat_h, MID,
                stroke=DEEP, stroke_pt=2.0, radius=True, radius_adj=0.20)
    text_box(s, x=cx - chat_d/2, y=cy - chat_h/2 + 0.15, w=chat_d, h=0.5,
             text="LLM / Chat", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(s, x=cx - chat_d/2, y=cy - chat_h/2 + 0.75, w=chat_d, h=0.7,
             text="модель + интерфейс\n+ память диалога",
             size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Top: Orchestrator
    orch_w = 3.5
    orch_h = 1.0
    filled_rect(s, cx - orch_w/2, cy - chat_h/2 - orch_h - 0.5, orch_w, orch_h,
                GOLD, stroke=DEEP, stroke_pt=1.5, radius=True, radius_adj=0.15)
    text_box(s, x=cx - orch_w/2, y=cy - chat_h/2 - orch_h - 0.45, w=orch_w, h=0.4,
             text="ОРКЕСТРАТОР", size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=cx - orch_w/2, y=cy - chat_h/2 - orch_h - 0.05, w=orch_w, h=0.5,
             text="планирование • decision loop",
             size=10, italic=True, color=DEEP, align=PP_ALIGN.CENTER)
    # Left: Memory
    mem_w = 2.6
    mem_h = 1.4
    mem_x = cx - chat_d/2 - mem_w - 0.6
    mem_y = cy - mem_h/2
    ocean_box(s, mem_x, mem_y, mem_w, mem_h, fill=LIGHT, stroke=DEEP)
    text_box(s, x=mem_x, y=mem_y + 0.15, w=mem_w, h=0.5,
             text="ПАМЯТЬ", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(s, x=mem_x, y=mem_y + 0.65, w=mem_w, h=0.7,
             text="vector DB,\nфайлы, логи",
             size=10, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Right: Tools
    tools_w = 2.6
    tools_h = 1.4
    tools_x = cx + chat_d/2 + 0.6
    tools_y = cy - tools_h/2
    ocean_box(s, tools_x, tools_y, tools_w, tools_h, fill=TEAL, stroke=DEEP)
    text_box(s, x=tools_x, y=tools_y + 0.15, w=tools_w, h=0.5,
             text="ИНСТРУМЕНТЫ", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(s, x=tools_x, y=tools_y + 0.65, w=tools_w, h=0.7,
             text="API, файлы,\ncode, search",
             size=10, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Connecting lines
    # Orchestrator → LLM (vertical)
    filled_rect(s, cx - 0.03, cy - chat_h/2 - 0.50, 0.06, 0.45, DEEP)
    # LLM ↔ Memory
    filled_rect(s, mem_x + mem_w + 0.05, cy - 0.03, 0.55, 0.06, DEEP)
    # LLM ↔ Tools
    filled_rect(s, cx + chat_d/2 + 0.03, cy - 0.03, 0.55, 0.06, DEEP)
    # Bottom: decision loop labels
    loop_y = 6.20
    text_box(s, x=0.55, y=loop_y, w=12.25, h=0.4,
             text="Цикл (ReAct): plan → act → observe → reflect → continue / stop",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=0.55, y=loop_y + 0.45, w=12.25, h=0.35,
             text="Yao et al. 2022 (arXiv:2210.03629)",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    """Agent: 200 PDF case + 5 levels of autonomy."""
    s = blank(p)
    slide_title(s, "Агент за работой: кейс 200 PDF + 5 уровней автономии.", size=26)
    # Left: case card
    cx_, cy_, cw_, ch_ = 0.55, 1.95, 5.4, 4.7
    ocean_box(s, cx_, cy_, cw_, ch_)
    text_box(s, x=cx_ + 0.30, y=cy_ + 0.25, w=cw_ - 0.6, h=0.4,
             text="Кейс — типовой для агента", size=14, bold=True, color=TEAL)
    text_box(s, x=cx_ + 0.30, y=cy_ + 0.75, w=cw_ - 0.6, h=1.4,
             text="200 PDF-отчётов.\nИзвлечь дату, контрагента, сумму.\nСобрать сводную таблицу.",
             size=16, bold=True, color=DEEP, line_spacing=1.30)
    # Mock arrow PDFs → table
    text_box(s, x=cx_ + 0.30, y=cy_ + 2.4, w=cw_ - 0.6, h=2.1,
             text="Не модель — нет специализированной модели «возьми 200 произвольных PDF».\n\nНе чат — некомфортно копировать 200 файлов в окно.\n\nАгент — естественный выбор.",
             size=12, color=DEEP, line_spacing=1.45, italic=True)
    # Right: ladder of 5 levels
    lx, ly, lw, lh = cx_ + cw_ + 0.35, 1.95, 6.5, 4.7
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, x=lx + 0.25, y=ly + 0.20, w=lw - 0.5, h=0.4,
             text="5 уровней автономии (Feng/McDonald/Zhang, 2025)",
             size=13, bold=True, color=DEEP)
    levels = [
        ("1. Operator", "пользователь на каждом шаге", "Claude Code «approve each»", LIGHT),
        ("2. Collaborator", "пара, ролями перетекая", "Cursor парное программирование", LIGHT),
        ("3. Consultant", "цель + план + правки", "Devin фиксит баг по тикету", MID),
        ("4. Approver", "agent действует, gate", "agent собирает PR, ждёт review", MID),
        ("5. Observer", "полная автономия", "AutoGPT на ночь", GOLD),
    ]
    rh = 0.72
    rt = ly + 0.80
    for i, (name, role, ex, color) in enumerate(levels):
        ry = rt + (4 - i) * rh  # bottom-up ladder visually
        is_gold = (color == GOLD)
        filled_rect(s, lx + 0.25, ry, lw - 0.5, rh - 0.05, color,
                    stroke=DEEP if is_gold else None, stroke_pt=1.5 if is_gold else 0.0,
                    radius=True, radius_adj=0.15)
        text_box(s, x=lx + 0.40, y=ry + 0.05, w=2.4, h=0.35, text=name,
                 size=13, bold=True, color=DEEP if is_gold else WHITE)
        text_box(s, x=lx + 0.40, y=ry + 0.36, w=lw - 0.8, h=0.30, text=role,
                 size=10, italic=True, color=DEEP if is_gold else WHITE)
        text_box(s, x=lx + 2.85, y=ry + 0.05, w=lw - 3.2, h=0.35, text=ex,
                 size=10, color=DEEP if is_gold else WHITE, align=PP_ALIGN.RIGHT)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    """Applications product UX — Translate metrics + 6 logos grid."""
    s = blank(p)
    slide_title(s, "Приложение = AI, упакованный в продуктовый интерфейс.", size=28)
    # Top: Translate metrics
    mt_x, mt_y, mt_w, mt_h = 0.55, 1.85, 12.25, 1.4
    ocean_box(s, mt_x, mt_y, mt_w, mt_h, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=mt_x + 0.3, y=mt_y + 0.20, w=mt_w - 0.6, h=0.4,
             text="Google Translate — масштаб 2026", size=14, bold=True, color=TEAL)
    text_runs(s, mt_x + 0.3, mt_y + 0.65, mt_w - 0.6, 0.7, [
        {"text": "1+ миллиард ", "size": 26, "bold": True, "color": DEEP},
        {"text": "уникальных пользователей в месяц  ·  ", "size": 14, "color": DEEP},
        {"text": "1 триллион ", "size": 26, "bold": True, "color": GOLD},
        {"text": "переведённых слов / месяц", "size": 14, "color": DEEP},
    ], line_spacing=1.0)
    text_box(s, x=mt_x + 0.3, y=mt_y + mt_h - 0.35, w=mt_w - 0.6, h=0.3,
             text="across Google Translate, Search, Lens и Circle to Search (Google Blog, апрель 2026)",
             size=10, italic=True, color=LIGHT)
    # 6 logo grid
    logos = [
        ("logo-googletranslate.png", "Google Translate", "нейронная трансляция"),
        ("logo-notion.png", "Notion AI", "GPT-4/Claude в кнопках"),
        ("logo-yandex.png", "ЯндексGPT в Поиске", "AI-краткий ответ"),
        ("logo-grammarly.png", "Grammarly", "NLP + LLM подсказки"),
        ("logo-yandex.png", "Яндекс.Карты", "ML маршрутизация"),
        ("logo-adobefirefly.png", "Adobe Firefly", "диффузия в Photoshop"),
    ]
    grid_y = 3.55
    cell_w = 4.05
    cell_h = 1.4
    grid_x = 0.55
    cell_gap = 0.10
    for i, (icon, name, role) in enumerate(logos):
        col = i % 3
        row = i // 3
        x = grid_x + col * (cell_w + cell_gap)
        y = grid_y + row * (cell_h + cell_gap)
        ocean_box(s, x, y, cell_w, cell_h)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + 0.20, y=y + 0.30, w=0.80, h=0.80)
        text_box(s, x=x + 1.15, y=y + 0.25, w=cell_w - 1.3, h=0.4, text=name,
                 size=13, bold=True, color=DEEP)
        text_box(s, x=x + 1.15, y=y + 0.70, w=cell_w - 1.3, h=0.55, text=role,
                 size=10, italic=True, color=SLATE, line_spacing=1.30)
    gold_callout(s, 0.55, 6.55, 12.25, 0.50,
                 "AI как функция, а не продукт. Большинство студентов уже ежедневно пользуются всеми шестью.",
                 size=12)
    speaker_notes(s, load_notes("s20"))


def build_s21(p):
    """Checklist: 2 questions + quadrant 2x2."""
    s = blank(p)
    slide_title(s, "Чек-лист «Какой тип AI выбрать»: 2 вопроса + квадрант 2×2.", size=24)
    # Top: 2 questions side by side — taller boxes
    q_y, q_h = 1.85, 1.30
    q1_x, q1_w = 0.55, 6.05
    q2_x, q2_w = q1_x + q1_w + 0.15, 6.05
    # Q1
    filled_rect(s, q1_x, q_y, q1_w, q_h, MID, radius=True, radius_adj=0.12)
    text_box(s, x=q1_x + 0.25, y=q_y + 0.15, w=q1_w - 0.5, h=0.35,
             text="ВОПРОС 1", size=12, bold=True, color=GOLD)
    text_box(s, x=q1_x + 0.25, y=q_y + 0.55, w=q1_w - 0.5, h=0.70,
             text="Нужно ли взаимодействие с пользователем?",
             size=14, bold=True, color=WHITE, line_spacing=1.25)
    # Q2
    filled_rect(s, q2_x, q_y, q2_w, q_h, MID, radius=True, radius_adj=0.12)
    text_box(s, x=q2_x + 0.25, y=q_y + 0.15, w=q2_w - 0.5, h=0.35,
             text="ВОПРОС 2", size=12, bold=True, color=GOLD)
    text_box(s, x=q2_x + 0.25, y=q_y + 0.55, w=q2_w - 0.5, h=0.70,
             text="Нужна ли самостоятельная работа с инструментами?",
             size=14, bold=True, color=WHITE, line_spacing=1.25)
    # Bottom: 2x2 quadrant — large; pushed down
    quad_x, quad_y = 1.7, 3.55
    quad_w, quad_h = 9.3, 3.05
    filled_rect(s, quad_x, quad_y, quad_w, quad_h, WHITE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.04)
    filled_rect(s, quad_x + quad_w / 2 - 0.005, quad_y, 0.01, quad_h, SOFT_GREY)
    filled_rect(s, quad_x, quad_y + quad_h / 2 - 0.005, quad_w, 0.01, SOFT_GREY)
    # Axis labels (column/row headers) — Q2 just above quadrant; Q1 inside the quadrant left edge per Fix-8
    text_box(s, x=quad_x + 0.10, y=quad_y - 0.30, w=quad_w / 2 - 0.20, h=0.28,
             text="Q2 = Нет (без инструментов)", size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(s, x=quad_x + quad_w / 2 + 0.10, y=quad_y - 0.30, w=quad_w / 2 - 0.20, h=0.28,
             text="Q2 = Да (с инструментами)", size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # Q1 labels — placed INSIDE the quadrant near the LEFT edge
    text_box(s, x=quad_x + 0.10, y=quad_y + 0.05, w=1.5, h=0.28,
             text="Q1 = Да ↑", size=10, bold=True, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    text_box(s, x=quad_x + 0.10, y=quad_y + quad_h - 0.32, w=1.5, h=0.28,
             text="Q1 = Нет ↓", size=10, bold=True, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # Cell labels with worked examples
    cells = [
        # (col, row, big, sub, color, dot_label)
        (0, 0, "ЧАТ", "корп.чат для\nразбора ТЗ", LIGHT, ""),
        (1, 0, "АГЕНТ", "200 PDF →\nтаблица", GOLD, "← worked example"),
        (0, 1, "МОДЕЛЬ", "конвейерный\nдетектор", LIGHT, ""),
        (1, 1, "ПРИЛОЖЕНИЕ", "ETL с AI-\nклассификатором (автоматизация)", LIGHT, ""),
    ]
    cw_ = quad_w / 2
    ch_ = quad_h / 2
    for col, row, big, sub, color, dot in cells:
        cx = quad_x + col * cw_
        cy = quad_y + row * ch_
        is_gold = (color == GOLD)
        if is_gold:
            filled_rect(s, cx + 0.12, cy + 0.12, cw_ - 0.24, ch_ - 0.24, GOLD_TINT,
                        radius=True, radius_adj=0.08)
        text_box(s, x=cx + 0.20, y=cy + 0.18, w=cw_ - 0.4, h=0.55, text=big,
                 size=22, bold=True, color=DEEP if is_gold else color, align=PP_ALIGN.CENTER, line_spacing=1.10)
        text_box(s, x=cx + 0.20, y=cy + 0.85, w=cw_ - 0.4, h=0.65, text=sub,
                 size=12, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Bottom retrieval prompt
    gold_callout(s, 0.55, 6.65, 12.25, 0.45,
                 "Подумайте 30 секунд: какой угол квадранта — ваш AI-инструмент?",
                 size=12)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """Section 4 divider — boundaries."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=0.4, w=12.25, h=1.5, text="Раздел  4",
             size=110, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    text_box(s, x=0.55, y=2.55, w=12.25, h=1.2,
             text="Границы AI — ваша зона ответственности",
             size=36, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    # 3 numbered cards
    reasons = [
        ("1", "Решение встроить AI", "уже на вас"),
        ("2", "Ошибки AI", "системны и предсказуемы"),
        ("3", "Граница «не умеет»", "тоже ваша"),
    ]
    card_y = 4.0
    card_w = 3.7
    card_h = 1.85
    gap = 0.25
    start_x = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2.0
    for i, (num, big, sub) in enumerate(reasons):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        text_box(s, x=x + 0.20, y=card_y + 0.30, w=0.7, h=0.95, text=num,
                 size=44, bold=True, color=GOLD)
        text_box(s, x=x + 1.10, y=card_y + 0.40, w=card_w - 1.3, h=0.5, text=big,
                 size=15, bold=True, color=DEEP, line_spacing=1.20)
        text_box(s, x=x + 1.10, y=card_y + 0.95, w=card_w - 1.3, h=0.7, text=sub,
                 size=12, italic=True, color=SLATE, line_spacing=1.30)
    roadmap_bar(s, 4)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    """Consumer vs enterprise — 2 columns + Samsung anchor + EU AI Act."""
    s = blank(p)
    slide_title(s, "Consumer vs enterprise — куда уходят ваши данные.", size=26)
    # Two columns
    col_y, col_h = 1.95, 3.5
    col_w = 6.05
    # Left consumer
    cx_ = 0.55
    ocean_box(s, cx_, col_y, col_w, col_h, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, x=cx_ + 0.25, y=col_y + 0.20, w=col_w - 0.5, h=0.45,
             text="ПОТРЕБИТЕЛЬСКИЕ ТАРИФЫ",
             size=14, bold=True, color=GOLD)
    text_box(s, x=cx_ + 0.25, y=col_y + 0.65, w=col_w - 0.5, h=0.5,
             text="данные → обучение по умолчанию",
             size=15, bold=True, color=DEEP, line_spacing=1.20)
    bullets_l = [
        "ChatGPT Free / Plus — train by default",
        "Anthropic Claude (с сент. 2025) — opt-in, 5 лет хранение",
        "Gemini Free — train + human review, 3 года",
        "YandexGPT Free — стандартная политика",
    ]
    for i, b in enumerate(bullets_l):
        text_box(s, x=cx_ + 0.30, y=col_y + 1.4 + i * 0.45, w=col_w - 0.6, h=0.40,
                 text=f"•  {b}", size=11, color=DEEP)
    # Right enterprise
    ex_ = cx_ + col_w + 0.30
    ocean_box(s, ex_, col_y, col_w, col_h, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=ex_ + 0.25, y=col_y + 0.20, w=col_w - 0.5, h=0.45,
             text="ENTERPRISE / API",
             size=14, bold=True, color=TEAL)
    text_box(s, x=ex_ + 0.25, y=col_y + 0.65, w=col_w - 0.5, h=0.5,
             text="данные ≠ обучение",
             size=15, bold=True, color=DEEP, line_spacing=1.20)
    bullets_r = [
        "ChatGPT Enterprise / Business — no training",
        "OpenAI API (с марта 2023) — no training",
        "Anthropic for Business — Zero Data Retention доступен",
        "Google Workspace / Vertex AI — no training",
    ]
    for i, b in enumerate(bullets_r):
        text_box(s, x=ex_ + 0.30, y=col_y + 1.4 + i * 0.45, w=col_w - 0.6, h=0.40,
                 text=f"•  {b}", size=11, color=DEEP)
    # Bottom: Samsung anchor + EU fines
    bot_y = col_y + col_h + 0.20
    s_x, s_w = 0.55, 7.5
    ocean_box(s, s_x, bot_y, s_w, 1.30, fill=WHITE, stroke=GOLD, stroke_pt=2.0)
    text_box(s, x=s_x + 0.20, y=bot_y + 0.10, w=s_w - 0.4, h=0.4,
             text="Samsung 2023 — канонический инцидент", size=13, bold=True, color=GOLD)
    text_box(s, x=s_x + 0.20, y=bot_y + 0.55, w=s_w - 0.4, h=0.7,
             text="3 эпизода за месяц (март–апрель): код, транскрипт совещания, тестовые последовательности → попали в датасет OpenAI. Самсунг ввёл запрет внешнего GenAI.",
             size=11, color=DEEP, line_spacing=1.30)
    # EU AI Act
    eu_x = s_x + s_w + 0.30
    eu_w = SLIDE_W_IN - eu_x - 0.55
    filled_rect(s, eu_x, bot_y, eu_w, 1.30, MID, radius=True, radius_adj=0.10)
    text_box(s, x=eu_x + 0.20, y=bot_y + 0.10, w=eu_w - 0.4, h=0.4,
             text="EU AI Act — штрафы", size=13, bold=True, color=WHITE)
    text_box(s, x=eu_x + 0.20, y=bot_y + 0.55, w=eu_w - 0.4, h=0.35,
             text="до 15M € / 3% оборота", size=12, color=WHITE, bold=True)
    text_box(s, x=eu_x + 0.20, y=bot_y + 0.90, w=eu_w - 0.4, h=0.35,
             text="до 35M € / 7% — за prohibited", size=11, color=GOLD, bold=True)
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    """Hallucinations — fake DOI prompt + Vectara HHEM range + AI knows all."""
    s = blank(p)
    slide_title(s, "Галлюцинации: AI уверенно генерирует несуществующие DOI.", size=26)
    # Left: prompt + 3 fake DOIs
    px, py, pw, ph = 0.55, 1.95, 7.5, 4.5
    ocean_box(s, px, py, pw, ph)
    text_box(s, x=px + 0.25, y=py + 0.20, w=pw - 0.5, h=0.45,
             text="Промпт", size=13, bold=True, color=TEAL)
    filled_rect(s, px + 0.25, py + 0.65, pw - 0.5, 0.7, SURFACE, stroke=SOFT_GREY,
                stroke_pt=1.0, radius=True, radius_adj=0.08)
    text_box(s, x=px + 0.40, y=py + 0.78, w=pw - 0.8, h=0.5,
             text='«Назови три статьи 2023-2024 по теме "сейсмостойкость подземных трубопроводов" с авторами, журналом и DOI».',
             size=11, italic=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=px + 0.25, y=py + 1.55, w=pw - 0.5, h=0.4,
             text="Ответ AI (3 фейк-ссылки):",
             size=13, bold=True, color=GOLD)
    fakes = [
        ("Petrov A., Smith J. (2023).", "Seismic Resilience of Small-Diameter Pipelines.", "DOI: 10.1016/j.engfailanal.2023.107214 ✗"),
        ("Ivanov K. et al. (2024).", "Underground Infrastructure Earthquake Response.", "DOI: 10.1080/15732479.2024.2218450 ✗"),
        ("Chen L., Brown R. (2023).", "Microscale Pipe Vibration Analysis.", "DOI: 10.1007/s11069-023-06122-1 ✗"),
    ]
    for i, (auth, title, doi) in enumerate(fakes):
        ry = py + 2.0 + i * 0.78
        text_box(s, x=px + 0.40, y=ry, w=pw - 0.8, h=0.30, text=auth,
                 size=10, bold=True, color=DEEP)
        text_box(s, x=px + 0.40, y=ry + 0.25, w=pw - 0.8, h=0.30, text=title,
                 size=10, italic=True, color=DEEP)
        text_box(s, x=px + 0.40, y=ry + 0.50, w=pw - 0.8, h=0.30, text=doi,
                 size=10, color=GOLD, bold=True)
    # Right: Vectara HHEM band
    rx, ry_, rw, rh = px + pw + 0.35, 1.95, 4.4, 3.2
    ocean_box(s, rx, ry_, rw, rh, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=rx + 0.25, y=ry_ + 0.20, w=rw - 0.5, h=0.4,
             text="Vectara HHEM (2025-26)", size=13, bold=True, color=TEAL)
    text_box(s, x=rx + 0.25, y=ry_ + 0.65, w=rw - 0.5, h=0.4,
             text="процент галлюцинаций", size=11, italic=True, color=LIGHT)
    # Range bar
    text_box(s, x=rx + 0.25, y=ry_ + 1.25, w=rw - 0.5, h=0.45,
             text="< 1%", size=24, bold=True, color=TEAL)
    text_box(s, x=rx + 0.25, y=ry_ + 1.65, w=rw - 0.5, h=0.35,
             text="суммаризация (Gemini 2.0 Flash)",
             size=10, italic=True, color=SLATE)
    text_box(s, x=rx + 0.25, y=ry_ + 2.20, w=rw - 0.5, h=0.45,
             text="10–15%", size=24, bold=True, color=GOLD)
    text_box(s, x=rx + 0.25, y=ry_ + 2.60, w=rw - 0.5, h=0.35,
             text="reasoning (многошаговое)",
             size=10, italic=True, color=SLATE)
    # Anti-pattern callout below
    gold_callout(s, rx, ry_ + rh + 0.20, rw, 1.10,
                 "Анти-паттерн: «AI знает всё». Любой ответ AI по фактическому вопросу — гипотеза для проверки.",
                 size=12)
    speaker_notes(s, load_notes("s24"))


def build_s25(p):
    """Bias / sycophancy / shift — 3 cards + GPT-4o timeline."""
    s = blank(p)
    slide_title(s, "Bias / sycophancy / distribution shift — три проявления одной природы.", size=24)
    cards = [
        ("Bias", "lucide-scale-blue.png",
         "Модель повторяет перекосы датасета.",
         "Скрининг резюме обучен на исторических данных — дискриминирует, не «решая», а статистически."),
        ("Sycophancy", "lucide-smartphone-blue.png",
         "Модель учится у RLHF поддакивать.",
         "Соглашается с явно неверным, чрезмерно хвалит — пользователь не замечает потери критики."),
        ("Distribution shift", "lucide-trending-up-blue.png",
         "Данные периода — устаревают.",
         "Модель на коде 2023 в 2026 предложит устаревшую библиотеку без явного сбоя."),
    ]
    card_y = 1.85
    card_w = 4.05
    card_h = 3.4
    gap = 0.10
    start_x = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2.0
    colors = [LIGHT, MID, DEEP]
    for i, (name, icon, def_, ex) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        color = colors[i]
        ocean_box(s, x, card_y, card_w, card_h, stroke=color)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + 0.30, y=card_y + 0.30, w=0.65, h=0.65)
        text_box(s, x=x + 1.10, y=card_y + 0.35, w=card_w - 1.3, h=0.5, text=name,
                 size=20, bold=True, color=color)
        text_box(s, x=x + 0.30, y=card_y + 1.20, w=card_w - 0.6, h=0.7, text=def_,
                 size=13, bold=True, color=DEEP, line_spacing=1.30)
        text_box(s, x=x + 0.30, y=card_y + 2.00, w=card_w - 0.6, h=1.2, text=ex,
                 size=11, italic=True, color=SLATE, line_spacing=1.40)
    # GPT-4o timeline
    tl_y = 5.50
    tl_h = 1.0
    ocean_box(s, 0.55, tl_y, 12.25, tl_h, fill=WHITE, stroke=GOLD, stroke_pt=2.0)
    text_box(s, x=0.75, y=tl_y + 0.10, w=11.85, h=0.4,
             text="GPT-4o sycophancy — апрель 2025", size=13, bold=True, color=GOLD)
    text_runs(s, 0.75, tl_y + 0.50, 11.85, 0.4, [
        {"text": "25 апр", "size": 14, "bold": True, "color": MID},
        {"text": " — релиз обновления   →   ", "size": 12, "color": DEEP},
        {"text": "28 апр", "size": 14, "bold": True, "color": MID},
        {"text": " — начало rollback (Altman в Twitter тем же вечером)   →   ", "size": 12, "color": DEEP},
        {"text": "29 апр", "size": 14, "bold": True, "color": MID},
        {"text": " — postmortem", "size": 12, "color": DEEP},
    ])
    # Bottom takeaway
    text_box(s, x=0.55, y=6.65, w=12.25, h=0.4,
             text="Общая причина: модель отражает данные, на которых обучена.",
             size=13, italic=True, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s25"))


def build_s26(p):
    """4 speakers AGI table (renamed from old build_s27 in v3.1)."""
    s = blank(p)
    slide_title(s, "Прогнозы AGI — 4 спикера, 4 разных стимула.", size=26)
    # Table
    tx, ty, tw = 0.55, 1.95, 12.25
    rh_head = 0.5
    rh_row = 1.0
    # Header
    cols = [
        ("Спикер", 2.4),
        ("Аффилиация", 2.0),
        ("Прогноз AGI", 4.0),
        ("Материальный интерес", 3.85),
    ]
    # Header bg
    filled_rect(s, tx, ty, tw, rh_head, MID)
    cur_x = tx
    for label, w in cols:
        text_box(s, x=cur_x + 0.15, y=ty + 0.10, w=w - 0.3, h=rh_head - 0.2,
                 text=label, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w
    # Rows
    rows = [
        ("Sam Altman", "OpenAI",
         "«Знаем как построить AGI; начало superintelligence» (янв. 2026)",
         "$100B+ раунды; контекст IPO"),
        ("Dario Amodei", "Anthropic",
         "«AGI через 2–3 года; нобелевский уровень за 2 года» (Давос 2026)",
         "Конкуренция с OpenAI; раунд 2026"),
        ("Demis Hassabis", "Google\nDeepMind",
         "«50% к концу декады; нужны прорывы в continual learning» (TIME100, 2025)",
         "Лидер community; больше доверия при осторожной позиции"),
        ("Yann LeCun", "AMI Labs\n(экс-Meta)",
         "«LLM не приведут к AGI; нужны world models, JEPA»",
         "Раунд $1B (март 2026) на альтернативный путь"),
    ]
    for i, (sp, af, pr, st) in enumerate(rows):
        rt = ty + rh_head + i * rh_row
        bg = SURFACE if i % 2 == 0 else WHITE
        filled_rect(s, tx, rt, tw, rh_row, bg, stroke=SOFT_GREY, stroke_pt=0.5)
        cur_x = tx
        for j, (text, w) in enumerate(zip([sp, af, pr, st], [c[1] for c in cols])):
            is_speaker = (j == 0)
            text_box(s, x=cur_x + 0.15, y=rt + 0.15, w=w - 0.3, h=rh_row - 0.30,
                     text=text, size=11.5 if is_speaker else 10.5,
                     bold=is_speaker, color=DEEP if is_speaker else SLATE,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.30)
            cur_x += w
    # Bottom takeaway
    gold_callout(s, tx, 6.75, tw, 0.50,
                 "Ни один из 4 спикеров не занимает нейтрально-научной позиции.",
                 size=12)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """Section 5 divider — NEW in v3.1 per DoD §10 + reader-rendered feedback."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # «Раздел 5» on its own line, large outline, top
    text_box(s, x=0.55, y=0.5, w=12.25, h=1.5, text="Раздел  5",
             size=110, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Title at safe vertical band
    text_box(s, x=0.55, y=3.0, w=12.25, h=1.4,
             text="Что забрать домой",
             size=40, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    # Frame phrase
    text_box(s, x=0.55, y=4.50, w=12.25, h=0.8,
             text="Резюме · задание к семинару 1 · карта семестра · тизер лекции 2.",
             size=22, italic=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.20)
    # Roadmap bar with «Вы здесь — раздел 5 из 5»
    roadmap_bar(s, 5)
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """Summary + homework — 3 takeaway cards + gold homework (renamed from old build_s29 in v3.1)."""
    s = blank(p)
    slide_title(s, "Что мы прошли + задание к семинару 1.", size=28)
    takeaways = [
        ("AI — спектр, не монолит", "Тип задачи × модальность × тип реализации.\nГрамотное обсуждение начинается с явной классификации."),
        ("Выбор типа AI — навык", "2 диагностических вопроса + квадрант 2×2.\nИнструмент, который вы применяете на семинарах."),
        ("Целеполагание у человека", "Все классы ошибок требуют человеческого контура.\nГраница «AI / не-AI» — ваша инженерная зона."),
    ]
    card_y = 1.95
    card_w = 4.05
    card_h = 3.0
    gap = 0.10
    start_x = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2.0
    colors = [LIGHT, MID, DEEP]
    for i, (head, body) in enumerate(takeaways):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h, stroke=colors[i])
        text_box(s, x=x + 0.20, y=card_y + 0.20, w=0.7, h=0.85, text=str(i + 1),
                 size=44, bold=True, color=colors[i])
        text_box(s, x=x + 1.0, y=card_y + 0.30, w=card_w - 1.2, h=0.85, text=head,
                 size=15, bold=True, color=DEEP, line_spacing=1.20)
        text_box(s, x=x + 0.30, y=card_y + 1.40, w=card_w - 0.6, h=1.5, text=body,
                 size=11.5, color=DEEP, line_spacing=1.40)
    # Gold homework
    hw_y = 5.20
    filled_rect(s, 0.55, hw_y, 12.25, 1.65, GOLD_TINT, stroke=GOLD, stroke_pt=2.5,
                radius=True, radius_adj=0.10)
    text_box(s, x=0.85, y=hw_y + 0.15, w=11.65, h=0.5,
             text="К СЕМИНАРУ 1", size=14, bold=True, color=GOLD)
    text_box(s, x=0.85, y=hw_y + 0.55, w=11.65, h=0.6,
             text="Возьмите свой AI-инструмент → пропустите через 2-вопросный квадрант → защитите выбор перед группой.",
             size=16, bold=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=0.85, y=hw_y + 1.15, w=11.65, h=0.45,
             text="Одностраничный разбор любого формата (текст / схема / таблица). Тема семинара — «Какой тип AI выбрать».",
             size=11, italic=True, color=DEEP)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """Course roadmap — 17 lectures × 3 modules (renamed from old build_s30 in v3.1)."""
    s = blank(p)
    slide_title(s, "Карта семестра: 17 лекций × 3 модуля.", size=28)
    # Three modules horizontal
    modules = [
        ("Модуль 1", "Основы\n+ знакомые индустрии", LIGHT, 1, 8,
         ["1. Введение", "2. Большие модели", "3. Агенты, RAG, API",
          "4. ПО", "5. Финансы / ритейл", "6. Инж. проект.",
          "7. Медицина", "8. Креативные ◆РК1"]),
        ("Модуль 2", "Высокотехнологичное\nматериальное произв-во", MID, 9, 12,
         ["9. Авиакосмос", "10. Сельское хоз-во",
          "11. Производство", "12. Цифровые двойники ◆РК2"]),
        ("Модуль 3", "Инфокоммуникации,\nнаука, добыча, синтез", DEEP, 13, 17,
         ["13. Логистика", "14. Телеком + cybersec", "15. Наука",
          "16. Нефтегаз", "17. Синтез ◆РК3"]),
    ]
    mod_y = 1.85
    mod_h = 4.85
    total_lectures = 17
    bar_x = 0.55
    bar_w = SLIDE_W_IN - 2 * 0.55
    # Width per lecture
    lec_w = bar_w / total_lectures
    cur_x = bar_x
    for label, sub, color, l_start, l_end, lectures in modules:
        n = l_end - l_start + 1
        m_w = n * lec_w
        ocean_box(s, cur_x, mod_y, m_w - 0.05, mod_h, fill=WHITE, stroke=color, stroke_pt=2.0)
        # Header band
        filled_rect(s, cur_x, mod_y, m_w - 0.05, 0.85, color, radius=True, radius_adj=0.10)
        text_box(s, x=cur_x, y=mod_y + 0.08, w=m_w - 0.05, h=0.35, text=label,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text_box(s, x=cur_x + 0.10, y=mod_y + 0.40, w=m_w - 0.25, h=0.5, text=sub,
                 size=10, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.20)
        # Lectures list — extra space after lecture 1 for "Вы здесь" tag
        gap_after_first = 0.32
        for j, lec in enumerate(lectures):
            ly = mod_y + 1.05 + j * 0.43 + (gap_after_first if j > 0 and l_start == 1 else 0)
            is_now = (l_start + j == 1)
            text_box(s, x=cur_x + 0.15, y=ly, w=m_w - 0.3, h=0.40, text=lec,
                     size=10.5 if is_now else 10,
                     bold=is_now, color=DEEP if not is_now else GOLD,
                     line_spacing=1.20)
            if is_now:
                # "Вы здесь" inline next to «1. Введение»
                text_box(s, x=cur_x + 1.5, y=ly, w=m_w - 1.6, h=0.4,
                         text="←  Вы здесь", size=10, bold=True, italic=True,
                         color=GOLD)
        cur_x += m_w
    # Bottom note
    text_box(s, x=0.55, y=6.85, w=12.25, h=0.35,
             text="◆ — рубежные контроли (РК1 на С8, РК2 на С12, РК3 — итоговый — на С17).",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """Lecture 2 teaser — callback to YOLO + 4 concepts in Russian (renamed from old build_s31 in v3.1)."""
    s = blank(p)
    slide_title(s, "Лекция 2: «Как работают современные большие модели».", size=28)
    # Left: YOLO callback mini frame
    yc_x, yc_y, yc_w, yc_h = 0.55, 1.95, 5.4, 4.4
    ocean_box(s, yc_x, yc_y, yc_w, yc_h)
    text_box(s, x=yc_x + 0.25, y=yc_y + 0.20, w=yc_w - 0.5, h=0.4,
             text="Callback к началу лекции", size=12, italic=True, color=TEAL)
    if (ASSETS / "illustrations/s01-yolo-mock.png").exists():
        # Smaller frame
        img_w = yc_w - 0.5
        img_h = img_w * 720 / 1280
        add_image(s, ASSETS / "illustrations/s01-yolo-mock.png",
                  x=yc_x + 0.25, y=yc_y + 0.7, w=img_w, h=img_h)
    text_box(s, x=yc_x + 0.25, y=yc_y + yc_h - 1.6, w=yc_w - 0.5, h=1.4,
             text="Как именно эта модель за 30 миллисекунд узнаёт людей? Внутренности — на лекции 2.",
             size=13, italic=True, color=DEEP, line_spacing=1.35)
    # Right: 4 concept cards 2×2
    cx_, cy_, cw_, ch_ = yc_x + yc_w + 0.35, 1.95, 6.5, 4.4
    concepts = [
        ("lucide-file-text-blue.png", "Токены", "единицы, на которые модель режет текст"),
        ("lucide-network-blue.png", "Эмбеддинги", "(векторные представления)\nадреса в смысловом пространстве"),
        ("lucide-eye-blue.png", "Внимание", "(attention)\nна какие части входа смотреть"),
        ("lucide-zap-blue.png", "Температура", "случайность выбора\nследующего токена"),
    ]
    sub_w = (cw_ - 0.20) / 2
    sub_h = (ch_ - 0.20) / 2
    for i, (icon, name, sub) in enumerate(concepts):
        col = i % 2
        row = i // 2
        x = cx_ + col * (sub_w + 0.15) + 0.05
        y = cy_ + row * (sub_h + 0.15) + 0.05
        ocean_box(s, x, y, sub_w, sub_h)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + 0.25, y=y + 0.25, w=0.55, h=0.55)
        text_box(s, x=x + 0.95, y=y + 0.25, w=sub_w - 1.1, h=0.5, text=name,
                 size=16, bold=True, color=MID)
        text_box(s, x=x + 0.25, y=y + 0.95, w=sub_w - 0.5, h=sub_h - 1.1, text=sub,
                 size=11, italic=True, color=DEEP, line_spacing=1.35)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """Q&A — minimal (renamed from old build_s32 in v3.1)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=1.9, w=12.25, h=2.4, text="Q&A",
             size=140, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.0)
    text_box(s, x=0.55, y=5.4, w=12.25, h=0.7, text="Спасибо",
             size=36, color=MID, align=PP_ALIGN.CENTER, italic=True)
    # Contact at bottom right
    text_box(s, x=8.0, y=6.8, w=4.85, h=0.4,
             text="контакты лектора — заполняется перед лекцией",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    speaker_notes(s, load_notes("s31"))


# ============================================================
# Main
# ============================================================
BUILDERS = [
    build_s01, build_s02, build_s02a, build_s03, build_s04, build_s05a, build_s05b,
    build_s06, build_s07, build_s08, build_s09,
    build_s10, build_s11, build_s12, build_s13,
    build_s14, build_s15, build_s16, build_s17, build_s18, build_s19, build_s20, build_s21,
    build_s22, build_s23, build_s24, build_s25, build_s26,
    # v3.1: removed build_s26-old (ARC-AGI) and build_s28-old (Pearl);
    # added NEW build_s27 (section 5 divider); renumbered s27→s26, s29→s28, s30→s29, s31→s30, s32→s31.
    build_s27, build_s28, build_s29, build_s30, build_s31,
]


def main():
    p = setup_pres()
    for build in BUILDERS:
        build(p)
    p.save(str(OUT))
    print(f"Saved {OUT} with {len(BUILDERS)} slides.")


if __name__ == "__main__":
    main()
