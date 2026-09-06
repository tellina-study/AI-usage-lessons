"""
Lecture 2 EN v1.0 (issue #188) — full deck, 47 slides, English re-render.

Source-of-truth: deck.en.yaml + slides-en/*.md (EN sources, already
translated/QA'd — used verbatim, no re-translation here).
Structurally ported from build_lec02.py (RU v3.0) — same layout/shape/
color/coordinate calls, only visible strings replaced with EN equivalents.
Order (identical to RU): s01 s02 s02a s03 s04 s04b | s05a s05 s06 s08 s09
  s10 s11 | s12a s12 s13 s14 s15 s17 | s18a s18 s19 s21 s22 s20 s23 s25 |
  s26a s26 s27 s28 s29 s30 s31 s32 | s33a s33 s34 s36 s37 |
  s35a s35 s38 s39 s40 s41 s42.

EN-specific content notes:
- s05: 3 markup examples are cat / hyperparameter / tokenization
  (o200k_base) — slides-en's own EN examples, not a literal port of RU's
  cat/tokenization/клубника set.
- s11: RU stays as the cross-lingual comparison subject (content, not
  translated away) — chart shows EN ~0.25, RU ~0.5 (gold), ZH ~0.8,
  Python ~0.4 tokens/char.
- s18/s19/s20: EN sentence "The cat ate the mouse because it was hungry" —
  attention resolves "it"->"cat" via thematic-role plausibility (not RU's
  grammatical-gender trick); gold cell/bar is "it"->"cat" = 0.7.
- s15: EN phrases "How to configure SSL" / "Installing an HTTPS
  certificate" / "Deploying a React component" / "Building a React app" /
  "Borscht recipe" — exact values per slides-en/s15.
- No footer/page numbers on this EN render (bilingual publish-channel
  convention — see main()).

Pipeline: python-pptx direct (same canonical path as RU builder, per
notes/mcp-limitations.md [#71-1] — full rebuild each iteration).
Palette LOCKED: Ocean #21295C/#065A82/#1C7293 + teal #028090 + gold #F0AB00.
Canvas 13.333"x7.5" (16:9).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED, mirrored from deck.yaml v2.0.1) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
SLATE_PN  = RGBColor(0x5B, 0x66, 0x78)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered/assets"
SLIDES_DIR = ROOT / "slides-en"
OUT = ROOT / "rendered/lec-02-en.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ============================================================
# Helpers (ported from v1.8 builder; refs system dropped)
# ============================================================
def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall(A_NS + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, A_NS + "effectLst")


def set_fill_alpha(shp, opacity_pct):
    """Полупрозрачная заливка: добавляет <a:alpha> в solidFill/srgbClr
    (opacity_pct: 0..100, 100 = непрозрачный). v2.1 — постер s01."""
    spPr = shp._element.spPr
    sf = spPr.find(A_NS + "solidFill")
    if sf is None:
        return
    srgb = sf.find(A_NS + "srgbClr")
    if srgb is None:
        return
    for el in srgb.findall(A_NS + "alpha"):
        srgb.remove(el)
    etree.SubElement(srgb, A_NS + "alpha").set(
        "val", str(int(opacity_pct * 1000)))


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = line_spacing
    r = para.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    """runs: list of dicts {text, size, bold, italic, color, font, newpara,
    align, line_spacing, space_after_pt}."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            para = tf.add_paragraph()
            para.alignment = cfg.get("align", align)
            para.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before_pt"):
                para.space_before = Pt(cfg["space_before_pt"])
        r = para.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=14, bold=True, stroke_pt=1.2, font=FONT_BODY):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    """Prefer only one of w/h — both stretches non-proportionally
    (notes/mcp-limitations.md #73-render-1)."""
    if not Path(path).exists():
        return
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w))
    elif h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.5, h=1.0, w=12.3, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.18,
                align=PP_ALIGN.LEFT):
    text_box(slide, x, y, w, h, text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                radius=True, radius_adj=0.12)
    text_box(slide, x + 0.2, y + 0.06, w - 0.4, h - 0.12, text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.22)


def right_arrow(slide, x, y, w=0.6, h=0.4, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def down_arrow(slide, x, y, w=0.4, h=0.45, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def line_arrow(slide, x1, y1, x2, y2, *, color=MID, w_pt=2.0, dash=None):
    """Прямая линия-стрелка (tailEnd triangle) — для веерных стрелок s20 и
    пунктирных коннекторов s35."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w_pt)
    if dash is not None:
        conn.line.dash_style = dash
    ln = conn.line._get_or_add_ln()
    end = ln.makeelement(A_NS + "tailEnd",
                         {"type": "triangle", "w": "med", "len": "med"})
    ln.append(end)
    disable_shadow(conn)
    return conn


def plain_line(slide, x1, y1, x2, y2, *, color=LIGHT, w_pt=1.5, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w_pt)
    if dash is not None:
        conn.line.dash_style = dash
    disable_shadow(conn)
    return conn


def left_arrow(slide, x, y, w=0.6, h=0.4, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def outline_big_text(slide, x, y, w, h, text, *, size=200, stroke=GOLD,
                     stroke_pt=2.6, align=PP_ALIGN.LEFT):
    """Decorative outline numeral: white fill + gold stroke via a:ln XML."""
    tb = text_box(slide, x, y, w, h, text, size=size, bold=True,
                  color=WHITE, align=align, line_spacing=0.95)
    r = tb.text_frame.paragraphs[0].runs[0]
    rPr = r._r.get_or_add_rPr()
    ln = rPr.makeelement(A_NS + "ln", {"w": str(int(stroke_pt * 12700))})
    sf = etree.SubElement(ln, A_NS + "solidFill")
    clr = etree.SubElement(sf, A_NS + "srgbClr")
    clr.set("val", "F0AB00" if stroke == GOLD else str(stroke))
    rPr.insert(0, ln)
    return tb


def speaker_notes(slide, text):
    """Readable paragraphs: blank-line-separated blocks each become one
    notes paragraph; single newlines collapse to spaces."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    blocks = [b.strip() for b in re.split(r'\n\s*\n', text.strip()) if b.strip()]
    if not blocks:
        blocks = [""]
    for i, block in enumerate(blocks):
        one = re.sub(r'\s*\n\s*', ' ', block)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = one


def load_notes(slide_id):
    files = sorted(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\Z)', md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes).strip()
    # v3.0: strip-safe внутренние маркеры источников ([VFY-day-of] /
    # [FACT-CHECK: …]) не выпекаются в speaker notes pptx (pre-gate grep = 0)
    notes = re.sub(r'\s*\[(?:VFY-day-of|FACT-CHECK[^\]]*)\]', '', notes)
    return notes


def page_number(slide, n, *, color=SLATE_PN):
    tb = slide.shapes.add_textbox(Inches(12.15), Inches(7.16), Inches(1.1),
                                  Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    r = para.add_run(); r.text = str(n)
    r.font.name = FONT_BODY; r.font.size = Pt(10)
    r.font.italic = True; r.font.color.rgb = color
    return tb


# ============================================================
# v2.0 pipeline progress bar (dividers) — 7 стадий конвейера s04b.
# Активная стадия(и) gold. Без «вы здесь», без минут.
# ============================================================
PIPE_STAGES = ["Text", "Tokens", "Vectors", "LLM",
               "Distribution", "Token", "Text"]


def pipeline_bar(slide, active, *, y=6.62, bar_h=0.48, passed=frozenset(),
                 muted_gold=False):
    """active: int|set — gold-подсветка. v3.0: passed — уже пройденные
    стадии (teal-tint, briefs s18a/s26a); muted_gold=True — ВСЯ лента в
    приглушённом gold (s33a: раздел логически «над» конвейером, не
    отдельная стадия; визуально слабее полного gold s35a)."""
    if isinstance(active, int):
        active = {active}
    n = len(PIPE_STAGES)
    arrow_w = 0.24
    total_w = 12.3
    cell_w = (total_w - arrow_w * (n - 1)) / n
    start_x = (SLIDE_W_IN - total_w) / 2.0
    for i, label in enumerate(PIPE_STAGES):
        x = start_x + i * (cell_w + arrow_w)
        is_act = (not muted_gold) and (i in active)
        if muted_gold:
            fill, color = GOLD_TINT, DEEP
        elif is_act:
            fill, color = GOLD, DEEP
        elif i in passed:
            fill, color = TEAL_TINT, TEAL
        else:
            fill, color = SOFT_GREY, SLATE
        filled_rect(slide, x, y, cell_w, bar_h, fill, radius=True,
                    radius_adj=0.28)
        sz = 11 if len(label) < 12 else 9.5
        text_box(slide, x - 0.05, y + 0.06, cell_w + 0.1, bar_h - 0.12,
                 label, size=sz, bold=is_act, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            ax = x + cell_w + 0.03
            right_arrow(slide, ax, y + bar_h / 2 - 0.07, w=arrow_w - 0.06,
                        h=0.14, fill=LIGHT)


def section_divider(p, *, section_n, sub_title, frame_phrase, tag,
                    active_stage, notes_id, frame_bar=False,
                    frame_size=20, illus=None, illus_caption=None,
                    passed=frozenset(), bar_muted=False):
    """v3.1 divider (#183 round 3, owner-мандат: «в заголовке каждого
    раздела должен быть мем или интересная картинка!»): 2-колоночная
    композиция — слева текстовый блок (номер раздела + подзаголовок +
    frame_phrase + tag), справа реальное фото/мем в Ocean rounded box,
    ЗАМЕТНОЕ (~35-38% площади слайда), единый паттерн на всех 6
    дивайдерах. pipeline_bar на всю ширину внизу, без изменений."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # --- Left column: text (width ~6.35", x=0.55..6.9) ---
    text_col_w = 6.35
    text_box(s, 0.55, 0.75, text_col_w, 1.9, f"Section {section_n}",
             size=92, bold=True, color=GOLD,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.55, 2.62, text_col_w, 1.05, sub_title,
             size=33, bold=True, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.08)
    text_box(s, 0.55, 3.80, text_col_w, 1.35, f"“{frame_phrase}”",
             size=frame_size, italic=True, color=MID, align=PP_ALIGN.LEFT,
             line_spacing=1.2)
    text_box(s, 0.55, 5.30, text_col_w, 0.45, tag,
             size=16, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    # --- Right column: real photo/meme in an Ocean rounded box
    #     (~5.4x4.55" ~24.6 sq.in ~25% of slide area) ---
    if illus is not None:
        img_x, img_y = 7.25, 0.85
        img_w, img_h = 5.55, 4.55
        ocean_box(s, img_x, img_y, img_w, img_h, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.5)
        pad = 0.16
        cap_h = 0.52 if illus_caption else 0
        _place_image_contain(s, illus, img_x + pad, img_y + pad,
                             img_w - 2 * pad, img_h - 2 * pad - cap_h)
        if illus_caption:
            text_box(s, img_x + pad, img_y + img_h - cap_h - 0.04,
                     img_w - 2 * pad, cap_h, illus_caption, size=9.5,
                     italic=True, color=LIGHT, align=PP_ALIGN.CENTER,
                     line_spacing=1.05)
    pipeline_bar(s, active_stage, passed=passed, muted_gold=bar_muted)
    if frame_bar:
        total_w = 12.3
        fx = (SLIDE_W_IN - total_w) / 2.0 - 0.14
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(fx), Inches(6.62 - 0.12),
                                 Inches(total_w + 0.28), Inches(0.48 + 0.24))
        shp.fill.background()
        shp.line.color.rgb = GOLD; shp.line.width = Pt(2.2)
        try:
            shp.adjustments[0] = 0.35
        except Exception:
            pass
        disable_shadow(shp)
    speaker_notes(s, load_notes(notes_id))
    return s


def _place_image_contain(slide, path, x, y, w, h):
    """Вписывает изображение в box (x,y,w,h) с сохранением пропорций
    (contain — не crop, не stretch), центрируя по обеим осям. Обходит
    #73-render-1 (add_picture с обоими w/h стрейчит непропорционально)."""
    if not Path(path).exists():
        return
    from PIL import Image as PILImage
    with PILImage.open(str(path)) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # шире бокса относительно — вписываем по ширине
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_y = y
        draw_x = x + (w - draw_w) / 2
    add_image(slide, path, x=draw_x, y=draw_y, w=draw_w)


# ============================================================
# Раздел 0
# ============================================================
def build_s01(p):
    """v3.1 meme hook (#183 round 3 pattern, ported): real recognizable
    meme template "Well yes, but actually no" (imgflip, still frame from
    "The Pirates! Band of Misfits", Aardman) — our T=0 question in the
    template's empty top field, built-in caption answers literally;
    hero >=40% area. Overlay text via text_box (not baked into pixels) —
    same image file reused for EN."""
    s = blank(p)
    # Hero: real meme template 1600x1218 (ratio 1.31), empty top ~21% of
    # height — our question goes there as a text layer over the image.
    hero_w = 7.6
    hero_h = hero_w * 1218 / 1600     # ~5.79" (~44% of slide area)
    hx = (SLIDE_W_IN - hero_w) / 2
    hy = 0.55
    add_image(s, ASSETS / "web/well-yes-actually-no-template.jpg",
              x=hx, y=hy, w=hero_w)
    # Empty white band of the template ~ 0..21% of image height
    blank_band_h = hero_h * 0.205
    text_box(s, hx + 0.35, hy + 0.10, hero_w - 0.7, blank_band_h - 0.15,
             "temperature=0 means the answer is always identical, right?",
             size=25, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    # The meme's own built-in caption "Well yes, but actually no" already
    # answers the question literally — add a short gloss line below.
    text_box(s, 0.55, hy + hero_h + 0.10, 12.23, 0.5,
             "Yes, formally deterministic. But no — not in practice. "
             "Why — today's topic.",
             size=16, italic=True, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Cover: large "02" outline gold + title + subtitle + hero motif
    (4-stage pipeline icon on the right)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    # "02" — outline gold, decorative, right side
    outline_big_text(s, 9.15, 0.30, 4.1, 3.4, "02", size=230,
                     align=PP_ALIGN.CENTER)
    # Lecture tag
    text_box(s, 0.7, 1.05, 8.0, 0.5, "LECTURE 2", size=18, bold=True,
             color=TEAL)
    filled_rect(s, 0.72, 1.58, 0.7, 0.05, TEAL)
    # Title
    text_box(s, 0.7, 2.0, 8.6, 2.3,
             "How Modern Large Language\nModels Work",
             size=44, bold=True, color=DEEP, line_spacing=1.08)
    # Subtitle
    filled_rect(s, 0.7, 4.55, 0.05, 0.62, GOLD)
    text_box(s, 0.95, 4.55, 10.6, 0.75,
             "The inference pipeline — and the boundaries that change "
             "engineering decisions",
             size=20, italic=True, color=MID, line_spacing=1.22)
    # Hero motif: 4-stage pipeline icon (token / vector / attention /
    # distribution)
    icons = [("[ ]", "token"), ("0.21", "vector"), ("⇄", "attention"),
             ("%", "distribution")]
    pipe_y = 5.75
    cx = 0.95
    for i, (glyph, label) in enumerate(icons):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(pipe_y),
                                 Inches(0.78), Inches(0.78))
        shp.fill.solid()
        shp.fill.fore_color.rgb = MID if i != 2 else TEAL
        shp.line.fill.background(); disable_shadow(shp)
        tf = shp.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        r = para.add_run(); r.text = glyph
        r.font.name = FONT_MONO; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = WHITE
        text_box(s, cx - 0.45, pipe_y + 0.88, 1.7, 0.35, label,
                 size=11.5, color=LIGHT, align=PP_ALIGN.CENTER)
        if i < 3:
            right_arrow(s, cx + 0.92, pipe_y + 0.30, w=0.5, h=0.18,
                        fill=LIGHT)
        cx += 1.58
    speaker_notes(s, load_notes("s02"))


S02A_ROWS = [
    ("0", "Introduction", "the frame and the pipeline as a whole", True),
    ("1", "Tokenization", "how the model sees your text", False),
    ("2", "Embeddings", "the space of meaning and the boundary of similarity", False),
    ("3", "Attention Mechanism", "what matters right now: roles, caching, long context", False),
    ("4", "Sampling", "from distribution to token: temperature, determinism, invisible tokens", False),
    ("5", "Model Types and Sizes", "what models run on, multimodality, the 2026 landscape", False),
    ("6", "Wrap-up", "assembling the pipeline, recap of the mechanisms", False),
]


def build_s02a(p):
    """Lecture map: 7 horizontal row-cards (new Section 5 "Model Types
    and Sizes"); no M-chips, no minutes; active Section 0 — gold
    outline."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    slide_title(s, "Lecture map — 7 sections", size=28,
                align=PP_ALIGN.CENTER, y=0.35, h=0.62)
    row_w, row_h, gap = 12.1, 0.72, 0.10
    x0 = (SLIDE_W_IN - row_w) / 2
    y = 1.10
    for num, name, desc, active in S02A_ROWS:
        if active:
            ocean_box(s, x0, y, row_w, row_h, fill=WHITE, stroke=GOLD,
                      stroke_pt=2.5)
        else:
            ocean_box(s, x0, y, row_w, row_h, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.2)
        text_box(s, x0 + 0.25, y, 0.6, row_h, num, size=27, bold=True,
                 color=GOLD if active else LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        text_runs(s, x0 + 1.05, y, row_w - 1.4, row_h, [
            {"text": name, "size": 16, "bold": True, "color": DEEP},
            {"text": "  —  " + desc, "size": 13, "italic": True,
             "color": SLATE},
        ], anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + gap
    speaker_notes(s, load_notes("s02a"))


def build_s03(p):
    """Today's object of study — the "model" layer; nested layers from
    Lecture 1 (Model gold) + 2 rows Layer/Today."""
    s = blank(p)
    slide_title(s, "Today's object — the “model” layer from "
                   "Lecture 1's four layers", size=26)
    # Nested layers on the left, bottom-aligned (outermost = Application)
    base_x, bottom = 0.85, 6.75
    layers = [  # (label, w, h, fill, stroke, stroke_pt)
        ("Application", 5.5, 5.0, WHITE, LIGHT, 1.2),
        ("Agent", 4.6, 3.9, SURFACE, LIGHT, 1.2),
        ("Chat", 3.7, 2.8, TEAL_TINT, TEAL, 1.2),
        ("Model", 2.8, 1.7, GOLD_TINT, GOLD, 2.2),
    ]
    for label, w, h, fill, stroke, sp in layers:
        x = base_x + (5.5 - w) / 2
        y = bottom - h
        ocean_box(s, x, y, w, h, fill=fill, stroke=stroke, stroke_pt=sp)
        text_box(s, x, y + 0.06, w, 0.4, label, size=14,
                 bold=(label == "Model"),
                 color=DEEP if label == "Model" else LIGHT,
                 align=PP_ALIGN.CENTER)
    # Right side — 2 rows
    ocean_box(s, 6.8, 2.0, 6.0, 1.6, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.3)
    text_runs(s, 7.0, 2.2, 5.6, 1.2, [
        {"text": "The “model” layer: ", "size": 17, "bold": True,
         "color": MID},
        {"text": "stateless inference — input goes in, a prediction comes "
                 "out, no memory between calls.", "size": 15.5,
         "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 6.8, 4.0, 6.0, 1.6, GOLD_TINT, stroke=GOLD, stroke_pt=1.6,
                radius=True, radius_adj=0.10)
    text_runs(s, 7.0, 4.2, 5.6, 1.2, [
        {"text": "Today: ", "size": 17, "bold": True, "color": DEEP},
        {"text": "we take apart what happens inside that inference — and "
                 "where its design changes engineering decisions.",
         "size": 15.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """Lecture goal (gold on "important details") + 7 promise chips for
    the 7-section structure (row of 4 + row of 3, no M-codes)."""
    s = blank(p)
    slide_title(s, "Lecture goal", size=24, color=MID)
    ocean_box(s, 0.7, 1.25, 11.93, 2.05, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 1.05, 1.42, 11.25, 1.75, [
        {"text": "“Examine how a language model works — and get "
                 "into the ", "size": 22, "bold": True, "color": DEEP},
        {"text": "important details", "size": 22, "bold": True,
         "color": GOLD},
        {"text": " that change how you build prompts, agents, and "
                 "decisions.”", "size": 22, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    promises = [
        "why a fixed “strawberry” answer proves nothing",
        "why a role-based prompt genuinely changes the answer",
        "what a 1M-token window can actually do",
        "why T=0 doesn't guarantee identical answers",
        "what invisible reasoning tokens actually cost",
        "small model vs. giant — what criterion to choose by",
        "what to replace blind trust in benchmarks with",
    ]
    gap_x, gap_y = 0.24, 0.24
    card_h = 1.12
    # Ряд 1 — 4 карточки, ряд 2 — 3 карточки (центрированы)
    rows = [promises[:4], promises[4:]]
    y = 3.70
    for row_items in rows:
        n = len(row_items)
        card_w = (12.23 - gap_x * (n - 1)) / n if n == 4 else 3.93
        x = (SLIDE_W_IN - card_w * n - gap_x * (n - 1)) / 2
        for txt in row_items:
            ocean_box(s, x, y, card_w, card_h, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.2)
            filled_rect(s, x + 0.14, y + 0.17, 0.09, card_h - 0.34, GOLD)
            text_box(s, x + 0.38, y + 0.08, card_w - 0.55, card_h - 0.16,
                     txt, size=12.5, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.12)
            x += card_w + gap_x
        y += card_h + gap_y
    speaker_notes(s, load_notes("s04"))


def build_s04b(p):
    """Keystone: 7-stage inference pipeline + autoregression loop (tokens
    one at a time, each appended to input) + example + 4 subcards + gold
    callout."""
    s = blank(p)
    slide_title(s, "Data flow in an LLM — there and back again", size=26)
    stages = [  # (name, example, caption)
        ("Text", "“Hello”", "words"),
        ("Tokens", "[Hel][lo]", "vocabulary IDs"),
        ("Vectors", "vec₁, vec₂", "numbers"),
        ("LLM", "attention", "inference"),
        ("Distribution", "p(token | context)", "probabilities"),
        ("Token", "chosen", "choice"),
        ("Text", "answer", "back to text"),
    ]
    n = len(stages)
    arrow_w = 0.30
    total_w = 12.5
    cell_w = (total_w - arrow_w * (n - 1)) / n
    x0 = (SLIDE_W_IN - total_w) / 2
    y0, cell_h = 2.02, 1.35
    # Autoregression loop above the pipeline: [Token] (i=5) -> back to
    # [Tokens] (i=1); caption above the line
    x_from = x0 + 5 * (cell_w + arrow_w) + cell_w / 2
    x_to = x0 + 1 * (cell_w + arrow_w) + cell_w / 2
    loop_y = 1.74
    plain_line(s, x_from, y0, x_from, loop_y, color=TEAL, w_pt=2.2)
    plain_line(s, x_from, loop_y, x_to, loop_y, color=TEAL, w_pt=2.2)
    line_arrow(s, x_to, loop_y, x_to, y0 - 0.02, color=TEAL, w_pt=2.2)
    text_box(s, x_to + 0.25, loop_y - 0.36, x_from - x_to - 0.5, 0.32,
             "⟲ tokens are generated one at a time; each is appended "
             "to the input",
             size=12, italic=True, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER)
    for i, (name, ex, cap) in enumerate(stages):
        x = x0 + i * (cell_w + arrow_w)
        is_llm = (name == "LLM")
        ocean_box(s, x, y0, cell_w, cell_h,
                  fill=GOLD_TINT if is_llm else SURFACE,
                  stroke=GOLD if is_llm else LIGHT,
                  stroke_pt=2.2 if is_llm else 1.4)
        nm_sz = 13 if len(name) < 12 else 10.5
        text_box(s, x - 0.06, y0 + 0.12, cell_w + 0.12, 0.4, name,
                 size=nm_sz, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        ex_sz = 10.5 if len(ex) < 14 else 8.5
        text_box(s, x - 0.10, y0 + 0.62, cell_w + 0.20, 0.6, ex,
                 size=ex_sz, color=MID, align=PP_ALIGN.CENTER,
                 font=FONT_MONO, line_spacing=1.1)
        text_box(s, x - 0.15, y0 + cell_h + 0.08, cell_w + 0.30, 0.5, cap,
                 size=10.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        if i < n - 1:
            right_arrow(s, x + cell_w + 0.03, y0 + cell_h / 2 - 0.08,
                        w=arrow_w - 0.06, h=0.16, fill=MID)
    # 4 section subcards
    subs = [("Section 1", "Text → Tokens"),
            ("Section 2", "Tokens → Vectors"),
            ("Section 3", "LLM: attention"),
            ("Section 4", "Distribution → Token")]
    sub_w, sub_h, gap = 2.95, 0.95, 0.22
    sx0 = (SLIDE_W_IN - sub_w * 4 - gap * 3) / 2
    sy = 4.38
    for i, (nm, rng) in enumerate(subs):
        x = sx0 + i * (sub_w + gap)
        ocean_box(s, x, sy, sub_w, sub_h, fill=WHITE, stroke=TEAL,
                  stroke_pt=1.3)
        text_box(s, x, sy + 0.12, sub_w, 0.35, nm, size=13.5, bold=True,
                 color=TEAL, align=PP_ALIGN.CENTER)
        text_box(s, x, sy + 0.50, sub_w, 0.35, rng, size=12, color=DEEP,
                 align=PP_ALIGN.CENTER)
    gold_callout(s, 2.7, 5.78, 7.93, 0.75,
                 "Words only exist at the edges; inside — vectors.",
                 size=17, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s04b"))


# ============================================================
# Section 1 — Tokenization
# ============================================================
def build_s05a(p):
    """Surprised Pikachu — direct illustration of "the model confidently
    answers wrong."""
    section_divider(
        p, section_n=1, sub_title="Tokenization",
        frame_phrase="How the model sees your text",
        tag="3 case studies · 3 failures", active_stage=1, notes_id="s05a",
        passed={0},
        illus=ASSETS / "web/surprised-pikachu-tokenize-en.jpg")


def token_chips_runs(pairs):
    """[(text, color)] -> run dicts, monospace."""
    return [{"text": t, "size": 16, "bold": True, "color": c,
             "font": FONT_MONO} for t, c in pairs]


def build_s05(p):
    """3 markup examples + gold callout + caption."""
    s = blank(p)
    slide_title(s, "A token is an ID from the model's vocabulary — "
                   "not a letter, not a word", size=26)
    rows = [
        ([("cat", DEEP)], [("[cat]", MID)], "1 token"),
        ([("hyperparameter", DEEP)],
         [("[hyper]", MID), ("[param]", TEAL), ("[eter]", LIGHT)],
         "3 tokens"),
        ([("tokenization", DEEP)], [("[token]", MID), ("[ization]", TEAL)],
         "2 tokens (o200k_base)"),
    ]
    y = 1.75
    for src, toks, cnt in rows:
        ocean_box(s, 0.9, y, 11.5, 0.92, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.3)
        runs = token_chips_runs(src)
        runs.append({"text": "   →   ", "size": 16, "color": SLATE})
        runs += token_chips_runs(toks)
        runs.append({"text": "   →   ", "size": 16, "color": SLATE})
        runs.append({"text": cnt, "size": 17, "bold": True, "color": DEEP})
        text_runs(s, 1.25, y + 0.12, 10.9, 0.68, runs,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += 1.12
    gold_callout(s, 0.9, 5.35, 11.5, 0.75,
                 "“On average: 1 token ≈ 4 characters in English "
                 "(≈ 2 characters in Russian)”",
                 size=17, align=PP_ALIGN.CENTER)
    text_box(s, 0.9, 6.35, 11.5, 0.5,
             "Vocabulary and model are two separate artifacts: the "
             "vocabulary is built before the model is trained, by a "
             "separate algorithm on its own corpus.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """BPE: 2 columns corpus -> vocabulary + gold callout + caption."""
    s = blank(p)
    slide_title(s, "BPE — a compromise between an alphabet and a "
                   "full-word vocabulary", size=26)
    text_box(s, 0.55, 1.30, 12.3, 0.5,
             "Not all letters (too long) and not all words (unfamiliar "
             "ones fall through) — frequent subsequences.",
             size=16, italic=True, color=MID)
    # Left column — corpus (height per content, v2.0.2 item 7)
    col_y, col_h = 2.05, 2.45
    ocean_box(s, 1.0, col_y, 4.7, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 1.3, col_y + 0.18, 4.1, 0.45, "Training corpus",
             size=17, bold=True, color=MID)
    text_runs(s, 1.3, col_y + 0.75, 4.1, 2.0, [
        {"text": "low", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": "lower", "size": 16, "font": FONT_MONO, "color": DEEP,
         "newpara": True, "space_before_pt": 6},
        {"text": "newest", "size": 16, "font": FONT_MONO, "color": DEEP,
         "newpara": True, "space_before_pt": 6},
        {"text": "widest", "size": 16, "font": FONT_MONO, "color": DEEP,
         "newpara": True, "space_before_pt": 6},
    ])
    right_arrow(s, 6.05, col_y + col_h / 2 - 0.25, w=1.1, h=0.5, fill=MID)
    # Right column — BPE vocabulary
    ocean_box(s, 7.5, col_y, 4.7, col_h, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_box(s, 7.8, col_y + 0.18, 4.1, 0.45, "BPE vocabulary",
             size=17, bold=True, color=TEAL)
    text_runs(s, 7.8, col_y + 0.75, 4.1, 2.0, [
        {"text": "low", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "er", "size": 16, "font": FONT_MONO, "color": TEAL},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "new", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "est", "size": 16, "font": FONT_MONO, "color": TEAL},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "wid", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": "+ individual characters, frequent syllables, whole "
                 "words", "size": 12.5, "italic": True, "color": SLATE,
         "newpara": True, "space_before_pt": 12},
    ])
    gold_callout(s, 1.0, 4.95, 11.2, 0.80,
                 "The vocabulary is built once, before training; at "
                 "inference — a lookup of ready-made rules, not a "
                 "computation.", size=16,
                 align=PP_ALIGN.CENTER)
    text_box(s, 1.0, 6.05, 11.2, 0.5,
             "Different vendors cut the same text differently: Claude, "
             "GPT, Gemini all have their own vocabularies and rules.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s06"))


def build_s08(p):
    """Mechanism on the left (compact) + Expanding Brain meme (4 panels,
    own text over the template) + patch-race timeline GPT-5.2 → GPT-5.5 →
    GPT-5.6 + StrawberryBench + callout."""
    s = blank(p)
    slide_title(s, "The patch race: they fixed strawberry, then "
                   "cranberry — one word at a time", size=24)
    # Left — mechanism (compact, frees space for the meme below)
    ocean_box(s, 0.55, 1.62, 5.0, 1.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.74, 4.4, 0.36, "Mechanism — letter blindness",
             size=14, bold=True, color=MID)
    text_runs(s, 0.85, 2.10, 4.4, 0.7, [
        {"text": "strawberry", "size": 13.5, "font": FONT_MONO,
         "color": DEEP},
        {"text": " → ", "size": 13.5, "color": SLATE},
        {"text": "[st][raw][berry]", "size": 12.5, "font": FONT_MONO,
         "bold": True, "color": MID},
        {"text": "  The model sees ", "size": 13, "color": DEEP},
        {"text": "3 tokens", "size": 13, "bold": True, "color": DEEP},
        {"text": ", not 10 letters.", "size": 13, "color": DEEP},
    ], line_spacing=1.15)
    # Expanding Brain meme (imgflip template, own text over it via PIL) —
    # same patch race as the timeline on the right, as a recognizable
    # image.
    meme_h = 2.35
    meme_w = meme_h * 804 / 992
    meme_x = 0.55 + (5.0 - meme_w) / 2
    meme_y = 3.05
    ocean_box(s, meme_x - 0.10, meme_y - 0.10, meme_w + 0.20,
              meme_h + 0.20, fill=WHITE, stroke=TEAL, stroke_pt=1.3)
    add_image(s, ASSETS / "web/expanding-brain-strawberry-en.jpg",
              x=meme_x, y=meme_y, h=meme_h)
    # Right — patch-race timeline (4 cards along an axis)
    cards = [
        ("GPT-5.2 · Dec 2025", "strawberry ✗ — “two r's in strawberry”",
         False),
        ("GPT-5.5 · Apr 2026",
         "strawberry ✓  /  cranberry ✗ — “two r's” instead of three",
         False),
        ("GPT-5.6 · Jul 2026",
         "cranberry ✓ — another viral case patched", True),
        ("StrawberryBench",
         "847 questions, 7 difficulty levels — systematic testing "
         "instead of one viral question", False),
    ]
    y = 1.55
    plain_line(s, 6.05, 1.8, 6.05, 5.15, color=LIGHT, w_pt=2.0)
    for title, body, is_gold in cards:
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.95),
                                 Inches(y + 0.40), Inches(0.2), Inches(0.2))
        shp.fill.solid()
        shp.fill.fore_color.rgb = GOLD if is_gold else LIGHT
        shp.line.fill.background(); disable_shadow(shp)
        if is_gold:
            filled_rect(s, 6.4, y, 6.4, 0.92, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.10)
        else:
            ocean_box(s, 6.4, y, 6.4, 0.92, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.2)
        text_box(s, 6.65, y + 0.07, 5.95, 0.36, title, size=13.5, bold=True,
                 color=DEEP if is_gold else MID)
        text_box(s, 6.65, y + 0.42, 5.95, 0.48, body, size=12, color=DEEP,
                 line_spacing=1.1)
        y += 1.0
    gold_callout(s, 0.55, 5.62, 12.25, 0.85,
                 "“Jagged intelligence”: skill level is set by training "
                 "data, not by “general intelligence” — the same model "
                 "wins Olympiad gold and fails letter counting.",
                 size=14.5, align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.60, 12.25, 0.4, [
        {"text": "What to do: ", "size": 13.5, "bold": True, "color": MID},
        {"text": "test your own domain's “cranberry” — obscure cases "
                 "nobody hyped about; going viral ≠ having the skill.",
         "size": 13.5, "color": DEEP},
    ], align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """Numbers and code: motivation line + 2 columns + "What to do"
    block at the bottom."""
    s = blank(p)
    slide_title(s, "The tokenizer cuts by frequency, not by structure",
                size=26, y=0.42, h=0.6)
    text_box(s, 0.55, 1.08, 12.3, 0.42,
             "Numbers and code are the most common “non-text” "
             "inputs: the model's arithmetic and your token budget both "
             "depend on how they're cut.",
             size=14.5, italic=True, color=MID)
    col_y, col_h = 1.62, 3.70
    # Column "Numbers"
    ocean_box(s, 0.55, col_y, 6.0, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, col_y + 0.18, 5.4, 0.45, "Numbers", size=18,
             bold=True, color=MID)
    text_runs(s, 0.85, col_y + 0.80, 5.45, 2.4, [
        {"text": "1000000", "size": 17, "font": FONT_MONO, "color": DEEP},
        {"text": " → ", "size": 17, "color": SLATE},
        {"text": "[100]", "size": 17, "font": FONT_MONO, "bold": True,
         "color": MID},
        {"text": "[000]", "size": 17, "font": FONT_MONO, "bold": True,
         "color": TEAL},
        {"text": "[0]", "size": 17, "font": FONT_MONO, "bold": True,
         "color": LIGHT},
        {"text": "Chunks left-to-right ≠ place values, which are "
                 "read right-to-left: boundaries don't line up, and "
                 "column addition breaks.",
         "size": 14, "color": DEEP, "newpara": True, "space_before_pt": 12},
        {"text": "Right-to-left cutting improves arithmetic; "
                 "task-specific schemes give ", "size": 14, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
        {"text": "up to +33% accuracy", "size": 15, "bold": True,
         "color": GOLD},
        {"text": " over standard cutting.", "size": 14, "color": DEEP},
    ], line_spacing=1.2)
    # Column "Code"
    ocean_box(s, 6.8, col_y, 6.0, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 7.1, col_y + 0.18, 5.4, 0.45, "Code", size=18, bold=True,
             color=TEAL)
    text_runs(s, 7.1, col_y + 0.80, 5.45, 2.4, [
        {"text": "GPT-2: ", "size": 15, "bold": True, "color": DEEP},
        {"text": "16 tokens", "size": 15.5, "bold": True, "color": MID},
        {"text": " for a level-4 indentation.", "size": 15, "color": DEEP},
        {"text": "GPT-4: ", "size": 15, "bold": True, "color": DEEP,
         "newpara": True, "space_before_pt": 12},
        {"text": "groups whitespace — the vocabulary gets fixed, but not "
                 "for every task at once.", "size": 15, "color": DEEP},
    ], line_spacing=1.22)
    # Math Lady — confusion at "[123][456][78]?" from the tokenizer's
    # number-chunking; symbols only, no baked RU text, reused as-is.
    _place_image_contain(s, ASSETS / "web/mathlady-tokens.jpg",
                         7.3, 3.62, 5.0, 1.58)
    # 3 tips + group heading (v2.0.2 item 9)
    text_box(s, 0.55, 5.44, 4.0, 0.32, "What to do:", size=14, bold=True,
             color=MID)
    tips = ["Digit separators (“1,234,567”)",
            "Offload computation to a tool",
            "Consistent indentation"]
    tip_w, gap = 3.95, 0.2
    x0 = (SLIDE_W_IN - tip_w * 3 - gap * 2) / 2
    for i, t in enumerate(tips):
        x = x0 + i * (tip_w + gap)
        ocean_box(s, x, 5.78, tip_w, 0.75, fill=WHITE, stroke=TEAL,
                  stroke_pt=1.3)
        text_box(s, x + 0.15, 5.84, tip_w - 0.3, 0.62, t, size=13,
                 bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    text_box(s, 0.55, 6.66, 12.25, 0.55,
             "Off-the-shelf chat products already route counting to "
             "built-in tools (code interpreter) automatically; call a "
             "tool yourself for non-standard cases and your own apps "
             "built on the API.",
             size=12, italic=True, color=MID, align=PP_ALIGN.CENTER,
             line_spacing=1.15)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    """Story / mechanism + direct Magikarp (Pokemon) illustration,
    medium-sized + fact + "Practical impact" block (diagnosis +
    sanitization)."""
    s = blank(p)
    slide_title(s, "Roughly 4% of the vocabulary are glitch tokens",
                size=26, y=0.42, h=0.6)
    col_y, col_h = 1.18, 3.42
    # Left — story + mechanism + Magikarp illustration
    ocean_box(s, 0.55, col_y, 3.95, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.8, col_y + 0.14, 3.45, 0.4, "SolidGoldMagikarp (2023)",
             size=13, bold=True, color=MID, font=FONT_MONO)
    text_box(s, 0.8, col_y + 0.54, 3.45, 0.62,
             "A Reddit username that ended up in GPT's vocabulary — the "
             "model couldn't repeat it (the token's name is a reference "
             "to the Pokemon Magikarp).",
             size=10.5, color=DEEP, line_spacing=1.1)
    magi_h = 1.20
    magi_w = magi_h * 820 / 669
    magi_x = 0.8 + (3.45 - magi_w) / 2
    magi_y = col_y + 1.24
    add_image(s, ASSETS / "web/magikarp-clean.png",
              x=magi_x, y=magi_y, h=magi_h)
    text_runs(s, 0.8, magi_y + magi_h + 0.10, 3.45, 0.62, [
        {"text": "Mechanism: ", "size": 10, "bold": True, "color": TEAL},
        {"text": "vocabulary corpus ≠ model training corpus → the "
                 "token's embedding stays at its random initialization.",
         "size": 10, "color": DEEP},
    ], line_spacing=1.1)
    # Center — "Practical impact" (main block)
    ocean_box(s, 4.7, col_y, 3.95, col_h, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_box(s, 4.95, col_y + 0.16, 3.45, 0.45, "Practical impact",
             size=14.5, bold=True, color=TEAL)
    text_runs(s, 4.95, col_y + 0.62, 3.45, 2.7, [
        {"text": "• Parsing failures on exotic strings — rare "
                 "identifiers, obfuscated text, unusual Unicode",
         "size": 12.5, "color": DEEP},
        {"text": "• Production risks — logs, user-generated IDs, "
                 "arbitrary user input", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
        {"text": "• Unexplained off-topic answers with no code error",
         "size": 12.5, "color": DEEP, "newpara": True,
         "space_before_pt": 8},
    ], line_spacing=1.18)
    # Right — GlitchMiner fact
    ocean_box(s, 8.85, col_y, 3.95, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 9.1, col_y + 0.16, 3.45, 0.45, "GlitchMiner (AAAI 2026)",
             size=14.5, bold=True, color=MID)
    text_runs(s, 9.1, col_y + 0.62, 3.45, 2.7, [
        {"text": "roughly 4% of the vocabulary", "size": 16, "bold": True,
         "color": GOLD},
        {"text": " by one estimate;", "size": 13.5, "color": DEEP},
        {"text": "reproduced in the open Llama, Qwen, Gemma, Phi-3, "
                 "Mistral families.", "size": 13.5, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
    ], line_spacing=1.22)
    # "Practical impact" block — 2 cards
    text_box(s, 0.55, 4.76, 5.0, 0.35, "In practice:", size=14,
             bold=True, color=MID)
    ocean_box(s, 0.55, 5.14, 6.0, 1.15, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 0.8, 5.24, 5.55, 0.95,
             "Unexplained behavior on exotic strings (rare identifiers, "
             "obfuscated text, unusual Unicode)? Hypothesis: a glitch "
             "token. Test: replace the string with a placeholder.",
             size=12.5, color=DEEP, line_spacing=1.14)
    ocean_box(s, 6.8, 5.14, 6.0, 1.15, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 7.05, 5.24, 5.55, 0.95,
             "In products accepting arbitrary input — normalize and "
             "sanitize input before feeding the model.",
             size=12.5, color=DEEP, line_spacing=1.14,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.55, 6.48, 12.25, 0.4,
             "A systemic property of the pipeline, not a version bug — "
             "it doesn't get fixed by scale.",
             size=12.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """Language cost: QuickChart bar on the left + dynamics and callout
    on the right. Russian stays the cross-lingual comparison subject —
    content, not translated away."""
    s = blank(p)
    slide_title(s, "Russian text costs roughly 2× more than English",
                size=26)
    # Chart (980x560) — in an ocean box
    box_x, box_y, box_w, box_h = 0.55, 1.7, 7.3, 4.55
    ocean_box(s, box_x, box_y, box_w, box_h, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    img_w = 6.9
    add_image(s, ASSETS / "charts/s11-tokens-per-char-v2-en.png",
              x=box_x + 0.2, y=box_y + 0.28, w=img_w)
    text_box(s, box_x + 0.2, box_y + 0.28 + img_w * 560 / 980 + 0.06,
             img_w, 0.35, "GPT-family vocabularies",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right — dynamics (box extends down — bottom half holds the Always
    # Has Been meme "the Russian token costs more?")
    ocean_box(s, 8.15, 1.7, 4.65, 3.15, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.3)
    text_runs(s, 8.4, 1.9, 4.15, 1.1, [
        {"text": "Move to o200k_base:", "size": 15, "bold": True,
         "color": MID},
        {"text": "roughly ", "size": 15, "color": DEEP, "newpara": True,
         "space_before_pt": 8},
        {"text": "−35%", "size": 17, "bold": True, "color": DEEP},
        {"text": " for non-Latin languages — the gap narrows but doesn't "
                 "disappear.", "size": 15, "color": DEEP},
    ], line_spacing=1.25)
    _place_image_contain(s, ASSETS / "web/always-has-been-ru-cost-en.jpg",
                         8.35, 3.22, 4.25, 1.55)
    # "What to do" block: calibrate on your own language + when
    # translating to English pays off
    filled_rect(s, 8.15, 5.05, 4.65, 2.35, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.2, radius=True, radius_adj=0.08)
    text_runs(s, 8.38, 5.22, 4.2, 2.05, [
        {"text": "What to do:", "size": 14, "bold": True, "color": MID},
        {"text": "•  Calibrate any token-based limit for your own "
                 "language: retrieval chunks, max_tokens, context-window "
                 "budget.", "size": 12.5,
         "color": DEEP, "newpara": True, "space_before_pt": 6},
        {"text": "•  For batch processing of large volumes — consider "
                 "translating to English (≈2× cheaper); in interactive "
                 "work the difference isn't worth it.", "size": 12.5,
         "color": DEEP, "newpara": True,
         "space_before_pt": 6},
    ], line_spacing=1.18)
    speaker_notes(s, load_notes("s11"))


# ============================================================
# Section 2 — Embeddings
# ============================================================
def build_s12a(p):
    """Pam "They're the same picture" — a recognizable meme about
    "different words, one meaning," a direct nod to the section's topic
    (similarity/embeddings)."""
    section_divider(
        p, section_n=2, sub_title="Embeddings",
        frame_phrase="The space of meaning — and where similarity breaks "
                     "down",
        tag="4 case studies · 1 failure", active_stage=2, notes_id="s12a",
        passed={0, 1},
        illus=ASSETS / "web/pam-same-picture.jpg")


def build_s12(p):
    """Embedding = a lookup from the input table."""
    s = blank(p)
    slide_title(s, "Every token gets a vector from the model's input "
                   "table", size=26)
    # Main diagram
    ocean_box(s, 0.55, 1.95, 8.0, 3.05, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    # [cat] chip
    chip(s, 0.95, 3.05, 1.1, 0.62, "[cat]", fill=MID, color=WHITE, size=17,
         font=FONT_MONO)
    right_arrow(s, 2.2, 3.22, w=0.55, h=0.28, fill=MID)
    # Input table — mini 3x2 grid
    tbl_x, tbl_y = 2.95, 2.50
    text_box(s, tbl_x - 0.15, tbl_y - 0.42, 2.6, 0.38, "input table",
             size=13, bold=True, color=MID, align=PP_ALIGN.CENTER)
    for ri in range(3):
        for ci in range(2):
            fill = TEAL_TINT if ri == 1 else WHITE
            filled_rect(s, tbl_x + ci * 1.15, tbl_y + ri * 0.55, 1.1, 0.5,
                        fill, stroke=LIGHT, stroke_pt=1.0)
    text_box(s, tbl_x, tbl_y + 0.55, 1.1, 0.5, "[cat]", size=11,
             bold=True, color=DEEP, font=FONT_MONO, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, tbl_x + 1.15, tbl_y + 0.55, 1.1, 0.5, "→ vector", size=10,
             color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    right_arrow(s, 5.5, 3.22, w=0.55, h=0.28, fill=MID)
    text_box(s, 6.15, 2.92, 2.3, 0.9, "[ 0.21, −0.45,\n0.88, …, 0.13 ]",
             size=14, bold=True, color=DEEP, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)
    text_box(s, 0.85, 4.45, 7.4, 0.45,
             "learned during training along with the other weights; "
             "after training, the table is fixed",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Mini-callout of dimensions
    ocean_box(s, 8.85, 1.95, 3.95, 3.05, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 9.1, 2.15, 3.45, 0.4, "Dimensions", size=15, bold=True,
             color=TEAL)
    text_runs(s, 9.1, 2.60, 3.45, 2.1, [
        {"text": "text-embedding-3-small — 1536", "size": 13, "color": DEEP,
         "font": FONT_MONO},
        {"text": "text-embedding-3-large — 3072", "size": 13, "color": DEEP,
         "font": FONT_MONO, "newpara": True, "space_before_pt": 6},
        {"text": "flagship internal dimensions aren't published; "
                 "order of magnitude is thousands", "size": 12.5,
         "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 8},
    ], line_spacing=1.2)
    gold_callout(s, 0.55, 5.55, 12.25, 0.62,
                 "“Geometric closeness = semantic closeness”",
                 size=17, align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.28, 12.25, 0.5, [
        {"text": "What to do: ", "size": 13, "bold": True, "color": TEAL},
        {"text": "a typo or a different letter case is already a "
                 "different token and a different vector; normalize "
                 "input before embedding it.", "size": 13,
         "color": DEEP},
    ], align=PP_ALIGN.CENTER, line_spacing=1.15)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """Embeddings aren't just an inference internal but also a search
    tool; three lives in practice (3rd — gold, "your search/RAG") +
    callout about re-indexing."""
    s = blank(p)
    slide_title(s, "Embeddings aren't just an inference internal — "
                   "they're also a search tool", size=21.5, y=0.45, h=0.55)
    text_box(s, 0.55, 1.08, 12.3, 0.42,
             "When you build search or RAG, you're using the third life "
             "of the term “embedding” — a separate embedding model.",
             size=14.5, italic=True, color=MID)
    cards = [
        ("1", "The input lookup table", "inside inference",
         "Static: the vector for [cat] is the same in every sentence. "
         "A lookup by ID, before any context exists.", False),
        ("2", "The model's internal data representation",
         "inside inference",
         "Vectors the model recomputes as it reads context — after the "
         "attention layers. These are what carry the model's "
         "“understanding.”", False),
        ("3", "Vectors for search", "a standalone tool",
         "A vector for a whole text, from a separate embedding model — "
         "not the internals of your chat LLM. Its own product, its own "
         "training, its own leaderboards.", True),
    ]
    card_w, col_h, gap = 3.95, 3.62, 0.2
    x0 = (SLIDE_W_IN - card_w * 3 - gap * 2) / 2
    y0 = 1.68
    for i, (num, title, sub, body, is_gold) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        if is_gold:
            filled_rect(s, x, y0, card_w, col_h, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.06)
        else:
            ocean_box(s, x, y0, card_w, col_h, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.4)
        text_box(s, x + 0.25, y0 + 0.14, 0.8, 0.65, num, size=32, bold=True,
                 color=GOLD if is_gold else LIGHT)
        text_box(s, x + 0.25, y0 + 0.82, card_w - 0.5, 0.8, title,
                 size=16, bold=True, color=DEEP, line_spacing=1.1)
        text_box(s, x + 0.25, y0 + 1.62, card_w - 0.5, 0.35, sub,
                 size=11.5, italic=True, color=MID if is_gold else SLATE)
        text_box(s, x + 0.25, y0 + 2.02, card_w - 0.5, 1.5, body,
                 size=13, color=DEEP, line_spacing=1.2)
        if is_gold:
            text_box(s, x + 0.25, y0 + col_h - 0.44, card_w - 0.5, 0.38,
                     "this is your search/RAG", size=11.5, italic=True,
                     bold=True, color=MID)
    gold_callout(s, 0.55, 5.72, 12.25, 0.85,
                 "Updated the chat LLM → you do NOT need to re-index the "
                 "database: the index lives in the embedding model's "
                 "coordinate space.",
                 size=16, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """2D scatter with labeled feature axes + 3 fact cards."""
    s = blank(p)
    slide_title(s, "Tokens close in meaning sit near each other — "
                   "across hundreds to thousands of dimensions", size=24)
    # Left — scatter in an ocean box
    bx, by, bw, bh = 0.55, 1.65, 6.7, 5.0
    ocean_box(s, bx, by, bw, bh, fill=WHITE, stroke=LIGHT, stroke_pt=1.4)
    text_box(s, bx + 0.2, by + 0.12, 2.6, 0.35,
             "2D projection (PCA-style)", size=12.5, italic=True,
             color=SLATE)
    text_box(s, bx + 2.5, by + 0.12, bw - 2.7, 0.55,
             "each of the 1536+ axes is a learned feature; two are "
             "shown here",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.RIGHT,
             line_spacing=1.05)
    # Feature axes (labeled)
    ax_x, ax_y = 1.02, 6.02  # origin (bottom-left)
    line_arrow(s, ax_x, ax_y, bx + bw - 0.25, ax_y, color=LIGHT, w_pt=1.8)
    line_arrow(s, ax_x, ax_y, ax_x, by + 0.62, color=LIGHT, w_pt=1.8)
    text_box(s, ax_x + 0.1, ax_y + 0.10, bw - 0.9, 0.32,
             "axis ≈ feature: topic (web development ↔ cooking)",
             size=10.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_box(s, ax_x + 0.12, by + 0.62, 3.6, 0.32,
             "axis ≈ feature: infrastructure ↔ frontend",
             size=10.5, italic=True, color=LIGHT)
    # Points
    def dot(x, y, fill, r=0.17):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                 Inches(r * 2), Inches(r * 2))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = WHITE; shp.line.width = Pt(1.2)
        disable_shadow(shp)
    # Dashed cluster "clouds": SSL — top-left (infrastructure),
    # React — bottom-left (frontend); borscht — right (cooking)
    for (ex, ey, ew, eh) in [(1.35, 2.55, 2.95, 1.5),
                             (1.55, 4.25, 3.05, 1.5)]:
        ell = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ex), Inches(ey),
                                 Inches(ew), Inches(eh))
        ell.fill.background()
        ell.line.color.rgb = LIGHT; ell.line.width = Pt(1.2)
        ell.line.dash_style = 4  # dash
        disable_shadow(ell)
    # SSL points + labels
    dot(1.80, 2.85, MID)
    text_box(s, 2.15, 2.72, 2.2, 0.55, "How to configure SSL", size=11.5,
             bold=True, color=DEEP)
    dot(2.45, 3.42, MID)
    text_box(s, 2.90, 3.30, 2.3, 0.7, "Installing an\nHTTPS certificate",
             size=11.5, bold=True, color=DEEP, line_spacing=1.05)
    # React points + labels
    dot(2.00, 4.55, TEAL)
    text_box(s, 2.35, 4.42, 2.3, 0.55, "Deploying a React component",
             size=11.5, bold=True, color=DEEP)
    dot(2.70, 5.12, TEAL)
    text_box(s, 3.14, 5.00, 2.2, 0.7, "Building a\nReact app",
             size=11.5, bold=True, color=DEEP, line_spacing=1.05)
    # Outlier — borscht (right: cooking)
    dot(5.65, 3.95, GOLD)
    text_box(s, 5.30, 4.35, 1.8, 0.4, "Borscht recipe", size=11.5,
             bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, 5.05, 4.62, 2.3, 0.4, "outlier — a different area",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right — 3 fact cards
    facts = [
        ("Dimensionality", [
            {"text": "Public embedding models: ", "size": 13,
             "color": DEEP},
            {"text": "1536–3072", "size": 14.5, "bold": True, "color": GOLD},
            {"text": " dimensions; flagship internal dimensions are on "
                     "the order of thousands.",
             "size": 13, "color": DEEP}]),
        ("Training", [
            {"text": "Coordinates aren't hand-assigned: similar usage "
                     "contexts → close vectors.", "size": 13,
             "color": DEEP}]),
        ("Projection", [
            {"text": "You can only view the space through PCA/t-SNE — "
                     "the 2D picture loses some of the structure.",
             "size": 13, "color": DEEP}]),
    ]
    y = 1.65
    for title, runs in facts:
        ocean_box(s, 7.5, y, 5.3, 1.42, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.3)
        text_box(s, 7.75, y + 0.12, 4.8, 0.4, title, size=14.5, bold=True,
                 color=MID)
        text_runs(s, 7.75, y + 0.52, 4.8, 0.85, runs, line_spacing=1.18)
        y += 1.56
    gold_callout(s, 7.5, y + 0.08, 5.3, 1.0,
                 "What to do: closeness in this space is measurable as "
                 "distance — filtering and clustering without labels and "
                 "without an LLM is possible directly on the vectors, "
                 "cheaply.", size=12.5, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s14"))


# ---- s15 heatmap (native PowerPoint table, Ocean-scale fill) ----
S15_LABELS_ROW = ["How to configure SSL", "Installing an HTTPS certificate",
                  "Deploying a React component", "Building a React app",
                  "Borscht recipe"]
S15_LABELS_COL = ["SSL", "HTTPS", "React c.", "React a.", "Borscht"]
S15_VALS = [
    [1.00, 0.85, 0.18, 0.20, 0.08],
    [0.85, 1.00, 0.22, 0.19, 0.07],
    [0.18, 0.22, 1.00, 0.78, 0.12],
    [0.20, 0.19, 0.78, 1.00, 0.10],
    [0.08, 0.07, 0.12, 0.10, 1.00],
]


def _ocean_scale(v):
    """0..1 -> RGB on the SURFACE -> LIGHT -> DEEP scale."""
    stops = [(0.0, (0xF4, 0xF7, 0xFA)), (0.5, (0x1C, 0x72, 0x93)),
             (1.0, (0x21, 0x29, 0x5C))]
    for (t0, c0), (t1, c1) in zip(stops, stops[1:]):
        if v <= t1:
            f = (v - t0) / (t1 - t0)
            return RGBColor(*(int(a + (b - a) * f) for a, b in zip(c0, c1)))
    return RGBColor(*stops[-1][1])


def build_s15(p):
    """Similarity boundary: 5x5 heatmap (table with Ocean fill) +
    failure card + callout."""
    s = blank(p)
    slide_title(s, "High similarity means “about the same thing,” "
                   "not “with the same meaning”", size=24)
    # 6x6 table
    rows, cols = 6, 6
    tx, ty, tw, th = 0.55, 1.7, 7.15, 4.35
    gtbl = s.shapes.add_table(rows, cols, Inches(tx), Inches(ty),
                              Inches(tw), Inches(th))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(2.55)
    for ci in range(1, 6):
        tbl.columns[ci].width = Inches((tw - 2.55) / 5)
    tbl.rows[0].height = Inches(0.55)
    for ri in range(1, 6):
        tbl.rows[ri].height = Inches((th - 0.55) / 5)

    def cell_text(cell, txt, *, size=12, bold=False, color=DEEP,
                  align=PP_ALIGN.CENTER):
        cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        r = para.add_run(); r.text = txt
        r.font.name = FONT_BODY; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color

    # Corner cell
    tbl.cell(0, 0).fill.solid()
    tbl.cell(0, 0).fill.fore_color.rgb = WHITE
    cell_text(tbl.cell(0, 0), "cosine similarity", size=10.5, bold=True,
              color=SLATE, align=PP_ALIGN.LEFT)
    # Column headers
    for ci, lab in enumerate(S15_LABELS_COL):
        c = tbl.cell(0, ci + 1)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=11.5, bold=True, color=MID)
    # Rows
    for ri, lab in enumerate(S15_LABELS_ROW):
        c = tbl.cell(ri + 1, 0)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=11.5, bold=True, color=MID,
                  align=PP_ALIGN.LEFT)
        for ci, v in enumerate(S15_VALS[ri]):
            cell = tbl.cell(ri + 1, ci + 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ocean_scale(v)
            txt_color = WHITE if v >= 0.45 else DEEP
            hot = (v in (0.85, 0.78))
            cell_text(cell, f"{v:.2f}", size=12,
                      bold=(v >= 0.7), color=txt_color)
    # Right — failure card, gold
    filled_rect(s, 8.05, 1.7, 4.75, 1.68, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.2, radius=True, radius_adj=0.07)
    text_runs(s, 8.35, 1.84, 4.2, 1.45, [
        {"text": "“How to configure SSL”", "size": 15.5, "bold": True,
         "color": DEEP},
        {"text": "  ↔", "size": 15.5, "bold": True, "color": MID},
        {"text": "“How to disable SSL”", "size": 15.5, "bold": True,
         "color": DEEP, "newpara": True, "space_before_pt": 4},
        {"text": "Very high similarity — opposite practical meaning.",
         "size": 14, "color": DEEP, "newpara": True,
         "space_before_pt": 10},
    ], line_spacing=1.2)
    # Mini scale legend
    text_box(s, 8.05, 3.44, 4.75, 0.35, "scale: 0 — light · 1 — dark",
             size=11, italic=True, color=SLATE)
    # Spider-Man pointing at Spider-Man — "similar" != "about the same
    # thing" (similarity != relevance metaphor).
    ocean_box(s, 8.05, 3.82, 4.75, 2.42, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    _place_image_contain(s, ASSETS / "web/spiderman-similarity-en.jpg",
                         8.13, 3.89, 4.59, 2.28)
    gold_callout(s, 0.55, 6.3, 12.25, 0.8,
                 "Similarity is a candidate-generation signal; relevance "
                 "is a separate task: reranker, hybrid search, filters.",
                 size=15.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s15"))


def build_s17(p):
    """Section wrap-up: the round-trip path + 2 cards + callout +
    caption."""
    s = blank(p)
    slide_title(s, "Embeddings are the foundation of understanding: the "
                   "model works with vectors, not strings", size=23)
    # Left — vertical diagram
    bx, by, bw, bh = 0.55, 1.6, 3.5, 5.35
    ocean_box(s, bx, by, bw, bh, fill=SURFACE, stroke=LIGHT, stroke_pt=1.4)
    steps = ["words", "tokens", "vectors", "LLM", "vectors", "tokens",
             "words"]
    step_h, ar_h = 0.52, 0.15
    total = len(steps) * step_h + (len(steps) - 1) * ar_h
    yy = by + (bh - total) / 2
    for i, st in enumerate(steps):
        is_llm = (st == "LLM")
        filled_rect(s, bx + 0.55, yy, bw - 1.1, step_h,
                    GOLD if is_llm else (WHITE if i % 2 == 0 else TEAL_TINT),
                    stroke=GOLD if is_llm else LIGHT,
                    stroke_pt=1.6 if is_llm else 1.0,
                    radius=True, radius_adj=0.35)
        text_box(s, bx + 0.55, yy + 0.05, bw - 1.1, step_h - 0.1, st,
                 size=14, bold=True, color=DEEP if not is_llm else DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            down_arrow(s, bx + bw / 2 - 0.09, yy + step_h + 0.01, w=0.18,
                       h=ar_h - 0.02, fill=MID)
        yy += step_h + ar_h
    # Right — 2 cards
    cards = [
        ("Paraphrases and synonyms",
         "“How to set up SSL” and “Installing an HTTPS certificate” — "
         "close vectors → the model answers the same way; same with "
         "“car” and “automobile.”"),
        ("Cross-lingual proximity",
         "“клубника” and strawberry — close vectors → the answer is "
         "correct regardless of the query language."),
    ]
    y = 1.6
    for title, body in cards:
        ocean_box(s, 4.4, y, 8.4, 1.55, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.3)
        text_box(s, 4.7, y + 0.13, 7.8, 0.4, title, size=15, bold=True,
                 color=MID)
        text_box(s, 4.7, y + 0.55, 7.8, 0.9, body, size=13.5, color=DEEP,
                 line_spacing=1.2)
        y += 1.75
    gold_callout(s, 4.4, 5.15, 8.4, 0.72,
                 "Sentence-level semantic proximity is the basis of "
                 "“understanding” rephrasings.", size=14.5)
    text_runs(s, 4.4, 5.97, 8.4, 0.75, [
        {"text": "What to do: ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": "one embedding model and one index serve search, "
                 "clustering, and RAG all at once — but switching models "
                 "means re-indexing the entire store; choose it as an "
                 "infrastructure decision, not on the fly.",
         "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    text_box(s, 4.4, 6.85, 8.4, 0.45,
             "Choosing an embedding model for the task (MTEB, Matryoshka "
             "representations) — self-study material.",
             size=11.5, italic=True, color=LIGHT)
    speaker_notes(s, load_notes("s17"))


# ============================================================
# BATCH 2 — Section 3. Attention Mechanism
# ============================================================
def build_s18a(p):
    """Spotlight beam instead of an arXiv paper title page — "the beam
    lights up one thing, darkness around it" = a literal illustration of
    attention (what the model highlights from the context)."""
    section_divider(
        p, section_n=3, sub_title="Attention Mechanism",
        frame_phrase="How the model decides what to rely on in the "
                     "context — and what follows from that for roles, "
                     "caching, and long windows",
        tag="4 case studies · 2 failures", active_stage=3, notes_id="s18a",
        passed={0, 1, 2},
        frame_size=18,
        illus=ASSETS / "web/spotlight-clean.jpg")


# EN sentence: "The cat ate the mouse because it was hungry" — weight
# from "it" resolves to "cat" via thematic-role plausibility (the actor
# is the plausible bearer of hunger, not the mouse being eaten).
S18_TOKENS = ["The", "cat", "ate", "the", "mouse", "because", "it", "was",
              "hungry"]
S18_VALS = [
    [1.0, 0.3, 0.2, 0.1, 0.1, 0.1, 0.1, 0.0, 0.0],
    [0.4, 1.0, 0.3, 0.2, 0.2, 0.1, 0.1, 0.1, 0.1],
    [0.3, 0.4, 1.0, 0.2, 0.5, 0.1, 0.1, 0.1, 0.1],
    [0.1, 0.1, 0.1, 1.0, 0.3, 0.1, 0.1, 0.0, 0.0],
    [0.2, 0.2, 0.4, 0.3, 1.0, 0.1, 0.1, 0.0, 0.1],
    [0.1, 0.2, 0.2, 0.1, 0.1, 1.0, 0.3, 0.2, 0.2],
    [0.1, 0.7, 0.1, 0.1, 0.1, 0.2, 1.0, 0.3, 0.4],
    [0.1, 0.4, 0.1, 0.1, 0.1, 0.2, 0.4, 1.0, 0.5],
    [0.1, 0.4, 0.1, 0.1, 0.1, 0.2, 0.4, 0.5, 1.0],
]


def build_s18(p):
    """Attention matrix 9x9 (native table, Ocean scale, "it"->"cat"
    gold) + 3 properties + callout."""
    s = blank(p)
    slide_title(s, "Attention is a check of every token against every "
                   "other token", size=25, h=0.6)
    text_box(s, 0.55, 1.02, 12.3, 0.4,
             "Every token “looks at” every other token at once. "
             "At each step — N × N connections.",
             size=14, italic=True, color=MID)
    # 10x10 table (1 header + 9 tokens)
    n = len(S18_TOKENS)
    rows, cols = n + 1, n + 1
    tx, ty, tw, th = 0.55, 1.55, 7.5, 3.95
    first_col = 1.05
    gtbl = s.shapes.add_table(rows, cols, Inches(tx), Inches(ty),
                              Inches(tw), Inches(th))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(first_col)
    for ci in range(1, cols):
        tbl.columns[ci].width = Inches((tw - first_col) / n)
    tbl.rows[0].height = Inches(0.45)
    for ri in range(1, rows):
        tbl.rows[ri].height = Inches((th - 0.45) / n)

    def cell_text(cell, txt, *, size=9, bold=False, color=DEEP,
                  align=PP_ALIGN.CENTER):
        cell.margin_left = Inches(0.02); cell.margin_right = Inches(0.02)
        cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        r = para.add_run(); r.text = txt
        r.font.name = FONT_BODY; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color

    c00 = tbl.cell(0, 0)
    c00.fill.solid(); c00.fill.fore_color.rgb = WHITE
    cell_text(c00, "wt", size=8.5, bold=True, color=SLATE)
    for ci, lab in enumerate(S18_TOKENS):
        c = tbl.cell(0, ci + 1)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=8.5, bold=True, color=MID)
    for ri, lab in enumerate(S18_TOKENS):
        c = tbl.cell(ri + 1, 0)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=9, bold=True, color=MID,
                  align=PP_ALIGN.LEFT)
        for ci, v in enumerate(S18_VALS[ri]):
            cell = tbl.cell(ri + 1, ci + 1)
            cell.fill.solid()
            is_gold = (ri == 6 and ci == 1)   # "it" -> "cat"
            is_future = (ci > ri)  # upper triangle — future tokens
            if is_gold:
                cell.fill.fore_color.rgb = GOLD
                cell_text(cell, "0.7", size=10, bold=True, color=DEEP)
            elif is_future:
                # Grayed out: in the decoder a token doesn't see future
                # tokens
                cell.fill.fore_color.rgb = SOFT_GREY
                cell_text(cell, f"{v:.1f}", size=8,
                          color=SLATE)
            else:
                cell.fill.fore_color.rgb = _ocean_scale(v)
                cell_text(cell, f"{v:.1f}", size=9,
                          bold=(v >= 0.7),
                          color=WHITE if v >= 0.45 else DEEP)
    text_runs(s, 0.55, 5.62, 7.5, 0.8, [
        {"text": "In the row for “it”: ", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": "the largest weight lands on “cat,” not “mouse” — "
                 "resolved by semantic/thematic-role plausibility. A "
                 "statistical association learned from the corpus.",
         "size": 12.5, "color": DEEP},
        {"text": "In the decoder, a token sees only preceding tokens — "
                 "the full matrix is shown here for clarity.", "size": 11.5,
         "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 3},
    ], line_spacing=1.12)
    # Right — 3 properties
    props = [
        ("Dimensionality", [
            {"text": "N × N, where N is the context length. Doubling "
                     "the context ", "size": 12.5, "color": DEEP},
            {"text": "quadruples the attention compute", "size": 12.5,
             "bold": True, "color": DEEP},
            {"text": ".", "size": 12.5, "color": DEEP}]),
        ("At every step", [
            {"text": "The weight distribution is recomputed from "
                     "scratch at every generation step.", "size": 12.5,
             "color": DEEP}]),
        ("Multi-head", [
            {"text": "Each layer has dozens of parallel “heads” "
                     "(typically 32–128); each captures its own type of "
                     "relationship: grammar, semantics, long-range "
                     "dependencies.", "size": 12.5,
             "color": DEEP}]),
    ]
    y = 1.55
    for title, runs in props:
        ocean_box(s, 8.3, y, 4.5, 1.48, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.3)
        text_box(s, 8.55, y + 0.10, 4.0, 0.38, title, size=14, bold=True,
                 color=MID)
        text_runs(s, 8.55, y + 0.50, 4.0, 0.9, runs, line_spacing=1.15)
        y += 1.62
    gold_callout(s, 0.55, 6.45, 12.25, 0.70,
                 "What this affects: the weights determine whose Value "
                 "ends up in the current token's representation — and "
                 "directly shape the next prediction.",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    """Q/K/V on OUR sentence ("it was hungry") — worked example on the
    left, weight-distribution chart (leader "cat" gold) on the right +
    3 facts + flashlight metaphor."""
    s = blank(p)
    slide_title(s, "Attention returns a weight distribution: three "
                   "projections, Query / Key / Value",
                size=20, h=0.55, y=0.32)
    # Main thesis — right under the title (gold)
    gold_callout(s, 0.55, 1.06, 12.25, 0.6,
                 "Q is about the current step. K and V are about the "
                 "already-processed context.",
                 size=15.5, align=PP_ALIGN.CENTER)
    # Q/K/V row — three tiles
    qkv = [("Query", "“what I'm looking for right now”"),
           ("Key", "“what I offer”"),
           ("Value", "“what I hand over if I'm picked”")]
    x = 0.55
    for term, phrase in qkv:
        ocean_box(s, x, 1.84, 3.95, 0.58, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, x + 0.18, 1.84, 3.65, 0.58, [
            {"text": term, "size": 13.5, "bold": True, "color": MID,
             "font": FONT_MONO},
            {"text": " — " + phrase, "size": 12, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        x += 4.15
    # Left — worked example on our sentence
    ocean_box(s, 0.55, 2.62, 6.35, 3.05, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.3)
    text_runs(s, 0.85, 2.76, 5.8, 0.4, [
        {"text": "“The cat ate the mouse because ", "size": 14.5,
         "color": DEEP},
        {"text": "it", "size": 14.5, "bold": True, "color": MID},
        {"text": " was hungry”", "size": 14.5, "color": DEEP},
    ])
    text_runs(s, 0.85, 3.30, 5.8, 1.6, [
        {"text": "Q(“it”)", "size": 13, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": " = “looking for: who might have been hungry”",
         "size": 12.5, "color": DEEP},
        {"text": "K(“cat”)", "size": 13, "bold": True, "color": TEAL,
         "font": FONT_MONO, "newpara": True, "space_before_pt": 7},
        {"text": " = “I am an animate subject”", "size": 12.5,
         "color": DEEP},
        {"text": "V(“cat”)", "size": 13, "bold": True, "color": TEAL,
         "font": FONT_MONO, "newpara": True, "space_before_pt": 7},
        {"text": " = content that flows into “it”'s representation",
         "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    text_box(s, 0.85, 4.95, 5.8, 0.65,
             "A strong match between Q(“it”) and K(“cat”) gives a high "
             "weight → V(“cat”) determines the updated representation "
             "of the token “it.”",
             size=11.5, italic=True, color=MID, line_spacing=1.15)
    # Right — chart (9 tokens, same sentence split as in s18)
    ocean_box(s, 7.05, 2.62, 5.75, 3.05, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 7.3, 2.74, 5.25, 0.35,
             "Weight distribution — sum = 1, leader “cat”",
             size=12.5, bold=True, color=MID)
    add_image(s, ASSETS / "charts/s19-attention-weights-en.png",
              x=7.3, y=3.18, w=5.25)
    # 3 numbered facts — horizontal row
    triples = [
        ("1", "Input — all tokens in the context."),
        ("2", "Output — a weight distribution, sum = 1."),
        ("3", "Recomputed at every generation step."),
    ]
    x = 0.55
    for num, txt in triples:
        ocean_box(s, x, 5.85, 3.95, 0.62, fill=WHITE, stroke=TEAL,
                  stroke_pt=1.2)
        text_runs(s, x + 0.18, 5.85, 3.65, 0.62, [
            {"text": num + ".  ", "size": 13, "bold": True, "color": TEAL},
            {"text": txt, "size": 11.5, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
        x += 4.15
    # Flashlight metaphor — one caption line
    text_box(s, 0.55, 6.62, 12.3, 0.45,
             "Metaphor: a flashlight in a dark room — the beam points "
             "at relevant tokens, brightness = attention weight.",
             size=12.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    """Worked example "it"->"cat" (fan of arrows) + contrast with/without
    a role ("Explain the GIL") + research caveat (Zheng et al. — in
    notes) + gold callout."""
    s = blank(p)
    slide_title(s, "A role works through attention weight — but doesn't "
                   "raise factual quality", size=20, h=0.65, y=0.32)
    # Top box: worked example
    ocean_box(s, 0.55, 1.08, 12.25, 2.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.20, 5.5, 0.4, "Worked example: where does “it” "
             "look?", size=14, bold=True, color=MID)
    # Sentence tokens — separate boxes for the arrows
    seg = [("“The cat", 1.55, 1.65, True),
           ("ate the mouse,", 3.30, 2.15, False),
           ("because", 5.55, 1.35, False),
           ("it", 7.00, 0.45, True),
           ("was", 7.55, 0.75, False),
           ("hungry”", 8.40, 1.35, False)]
    for txt, x, w, bold in seg:
        text_box(s, x, 1.66, w, 0.45, txt, size=19, bold=bold,
                 color=DEEP if not bold else MID, line_spacing=1.0)
    # Arrows from "it" (x~7.20) to targets: "cat" (gold), "was", "hungry"
    line_arrow(s, 7.20, 2.13, 2.10, 2.50, color=GOLD, w_pt=4.0)
    line_arrow(s, 7.35, 2.13, 7.80, 2.50, color=MID, w_pt=2.2)
    line_arrow(s, 7.45, 2.13, 9.00, 2.50, color=LIGHT, w_pt=1.2)
    text_box(s, 2.10, 2.56, 2.4, 0.32, "main weight", size=11.5, bold=True,
             color=GOLD)
    text_box(s, 0.85, 2.90, 6.8, 0.42,
             "Simplification: an aggregate of hundreds of connections "
             "across dozens of layers — the model reproduces "
             "correlations of usage, not grammatical parsing.",
             size=10.5, italic=True, color=SLATE,
             line_spacing=1.1)
    text_runs(s, 7.85, 2.68, 4.7, 0.62, [
        {"text": "Think for 30 seconds: ", "size": 11.5, "bold": True,
         "color": DEEP},
        {"text": "where will the weight from “it” go in “The server "
                 "crashed because it ran out of memory”?", "size": 11.5,
         "color": DEEP},
    ], line_spacing=1.1)
    # Contrast: no role / with role
    ocean_box(s, 0.55, 3.48, 5.7, 1.5, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 0.85, 3.60, 5.1, 0.38, "No role", size=15, bold=True,
             color=MID)
    text_runs(s, 0.85, 4.00, 5.15, 0.9, [
        {"text": "“Explain the GIL”", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": " → a neutral, generic answer.", "size": 12.5,
         "color": DEEP},
    ], line_spacing=1.18)
    right_arrow(s, 6.38, 4.05, w=0.6, h=0.34, fill=GOLD)
    ocean_box(s, 7.1, 3.48, 5.7, 1.5, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 7.4, 3.60, 5.1, 0.38, "With a role", size=15, bold=True,
             color=TEAL)
    text_runs(s, 7.4, 4.00, 5.15, 0.95, [
        {"text": "“You are an experienced Python developer", "size": 12.5,
         "bold": True, "color": TEAL},
        {"text": ". Explain the GIL” — role tokens get weight → shift "
                 "the distribution of subsequent tokens: more specific, "
                 "expert register.", "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    # Research caveat: Zheng et al. measured ONLY factual accuracy
    # (2410 questions, 162 personas) — split into 2 separate lines:
    # (1) exact quote with numbers, (2) course observation on
    # tone/style, not tied to Zheng et al.
    ocean_box(s, 0.55, 5.05, 12.25, 0.82, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 0.85, 5.11, 11.7, 0.36,
             "Zheng et al. (2024, EMNLP Findings; 2410 questions, 162 "
             "personas): a persona/role in the prompt does not raise "
             "factual accuracy — the effect of any specific role is "
             "unpredictable.",
             size=12.5, bold=True, color=DEEP, line_spacing=1.1)
    text_box(s, 0.85, 5.50, 11.7, 0.32,
             "Separately from that study — from course observations: a "
             "role noticeably changes the tone, style, and content "
             "selection of the answer.",
             size=11.5, italic=True, color=TEAL, line_spacing=1.05)
    gold_callout(s, 0.55, 6.0, 12.25, 0.95,
                 "A role is a tool for controlling style and focus, not "
                 "an “intelligence booster.” If you need an answer "
                 "grounded in your data, give it the data — not a third "
                 "adjective attached to the word “expert.”",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s20"))


def build_s21(p):
    """KV-cache: schema of K/V in cache + Q of a new token; prefill/decode
    phases; gold callout about long chats."""
    s = blank(p)
    slide_title(s, "K and V are cached — only Q is recomputed", size=26,
                h=0.6)
    text_runs(s, 0.55, 1.08, 12.3, 0.55, [
        {"text": "KV-cache: ", "size": 14.5, "bold": True, "color": MID},
        {"text": "the Key/Value of already-processed tokens are stored "
                 "in accelerator memory. At each step, only the Q of the "
                 "new token is computed — against the stored K/V.",
         "size": 14.5, "color": DEEP},
    ], line_spacing=1.15)
    # Schema
    ocean_box(s, 0.55, 1.75, 12.25, 1.85, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    xx = 1.0
    for i in range(4):
        lab = f"token {i+1}" if i < 3 else "…"
        chip(s, xx, 2.0, 1.15, 0.42, lab, fill=MID, color=WHITE, size=12)
        for j, kv in enumerate(("K", "V")):
            filled_rect(s, xx + 0.06 + j * 0.56, 2.52, 0.5, 0.42, TEAL_TINT,
                        stroke=TEAL, stroke_pt=1.2, radius=True,
                        radius_adj=0.25)
            text_box(s, xx + 0.06 + j * 0.56, 2.56, 0.5, 0.34, kv, size=13,
                     bold=True, color=TEAL, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
        xx += 1.45
    text_box(s, 1.0, 3.06, 5.6, 0.35,
             "already computed — sitting in cache, not recomputed",
             size=11, italic=True, color=TEAL)
    # Divider and new token
    plain_line(s, 7.15, 1.95, 7.15, 3.35, color=LIGHT, w_pt=1.2, dash=4)
    chip(s, 9.7, 2.0, 1.55, 0.42, "new token", fill=DEEP, color=WHITE,
         size=12)
    filled_rect(s, 10.2, 2.52, 0.5, 0.42, GOLD, stroke=None, radius=True,
                radius_adj=0.25)
    text_box(s, 10.2, 2.56, 0.5, 0.34, "Q", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    left_arrow(s, 7.45, 2.6, w=2.6, h=0.26, fill=GOLD)
    text_box(s, 7.45, 2.95, 2.7, 0.35, "checked against all K/V",
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)
    # Phases
    ocean_box(s, 0.55, 3.72, 5.9, 1.95, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 0.85, 3.84, 5.3, 0.4,
             "Phase 1 — prefill (processing the prompt)",
             size=14, bold=True, color=MID)
    text_runs(s, 0.85, 4.26, 5.35, 1.35, [
        {"text": "• All input tokens are known at once → K/V are "
                 "computed ", "size": 12.5, "color": DEEP},
        {"text": "in parallel", "size": 12.5, "bold": True, "color": DEEP},
        {"text": "• Bound by ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "compute power", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": "• Determines the ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "pause before the first character of the answer (TTFT)",
         "size": 12.5, "bold": True, "color": DEEP},
    ], line_spacing=1.14)
    right_arrow(s, 6.55, 4.55, w=0.45, h=0.28, fill=MID)
    ocean_box(s, 7.1, 3.72, 5.7, 1.95, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 7.4, 3.84, 5.1, 0.4,
             "Phase 2 — decode (generating the answer)",
             size=14, bold=True, color=TEAL)
    text_runs(s, 7.4, 4.26, 5.15, 1.35, [
        {"text": "• Strictly ", "size": 12.5, "color": DEEP},
        {"text": "sequential", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": ", token by token", "size": 12.5, "color": DEEP},
        {"text": "• Every step reads ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "the entire accumulated cache", "size": 12.5,
         "bold": True, "color": DEEP},
        {"text": " from memory", "size": 12.5, "color": DEEP},
        {"text": "• Bound by ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "memory bandwidth", "size": 12.5,
         "bold": True, "color": DEEP},
        {"text": " → “typing” speed", "size": 12.5, "color": DEEP},
    ], line_spacing=1.14)
    # Vendors: who caches automatically, who requires an explicit marker
    ocean_box(s, 0.55, 5.80, 12.25, 0.60, fill=SURFACE, stroke=TEAL,
              stroke_pt=1.2)
    text_runs(s, 0.85, 5.80, 11.7, 0.60, [
        {"text": "Cache by provider:  ", "size": 12.5, "bold": True,
         "color": TEAL},
        {"text": "OpenAI — automatic  ·  Gemini — implicit (automatic)  "
                 "·  DeepSeek — automatic, disk-based  ·  Anthropic — "
                 "explicit (cache_control)", "size": 12.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.52, 12.25, 0.75,
                 "A working KV-cache makes resubmitting history cheap "
                 "and fast. “Slow and expensive” is what happens when "
                 "the cache MISSES — not an inherent property of long "
                 "chats.",
                 size=13.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """Prompt caching: schema of THREE requests (write -> hit -> miss) +
    case bars on the left; exact-prefix stack on the right; gold
    callout."""
    s = blank(p)
    slide_title(s, "Prompt caching is a bet on reuse, not a discount",
                size=25, h=0.6)
    # Left — schema of three requests
    ocean_box(s, 0.55, 1.3, 6.35, 4.35, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.42, 5.8, 0.38, "Three requests in a row:",
             size=14.5, bold=True, color=MID)
    reqs = [
        ("Request 1", "the prefix is written to the cache — write "
         "costs 1.25–2× the rate", TEAL_TINT, TEAL, False),
        ("Request 2", "same prefix → cache hit — reads cost 0.1× "
         "(newest models — 0.025×)", GOLD_TINT, GOLD, True),
        ("Request 3", "a line was added at the front → prefix no "
         "longer matches → cache miss, full price", SOFT_GREY, SLATE,
         False),
    ]
    yy = 1.86
    for name, desc, fill, stroke, is_gold in reqs:
        filled_rect(s, 0.85, yy, 5.75, 0.62, fill, stroke=stroke,
                    stroke_pt=1.6 if is_gold else 1.2, radius=True,
                    radius_adj=0.14)
        text_runs(s, 1.02, yy + 0.03, 5.45, 0.56, [
            {"text": name + ":  ", "size": 12, "bold": True,
             "color": DEEP},
            {"text": desc, "size": 11.5, "color": DEEP,
             "bold": is_gold},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
        yy += 0.72
    # Case — 2 bars
    text_box(s, 0.85, 4.12, 5.8, 0.35,
             "Case: 50,000 document analyses per month", size=13,
             bold=True, color=DEEP)
    # Bars shortened (proportions kept) — Stonks (~1.65") sits on the
    # right as an emotional anchor for the -82% saving.
    filled_rect(s, 0.85, 4.50, 3.1, 0.32, MID)
    text_box(s, 4.05, 4.50, 1.1, 0.32, "$45,000", size=12, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.98, 4.52, 2.9, 0.28, "without cache", size=11,
             color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 0.85, 4.92, 0.55, 0.32, GOLD)
    text_runs(s, 1.5, 4.92, 3.6, 0.32, [
        {"text": "$8,000 with cache · ", "size": 12, "bold": True,
         "color": DEEP},
        {"text": "−82%", "size": 13.5, "bold": True, "color": GOLD},
    ], anchor=MSO_ANCHOR.MIDDLE)
    _place_image_contain(s, ASSETS / "web/stonks-template.png",
                         5.15, 4.35, 1.65, 1.25)
    # Right — exact prefix (illustrates request 3)
    ocean_box(s, 7.15, 1.3, 5.65, 4.35, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 7.4, 1.45, 5.15, 0.75, [
        {"text": "Condition: exact prefix match. ", "size": 13,
         "bold": True, "color": DEEP},
        {"text": "One changed token invalidates the cache for "
                 "everything after it.", "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    stack = [
        ("system prompt", TEAL_TINT, TEAL, "✓ cache"),
        ("instructions", TEAL_TINT, TEAL, "✓ cache"),
        ("documents", TEAL_TINT, TEAL, "✓ cache"),
        ("“Today is 2026-09-05 14:23” — dynamic", GOLD_TINT, GOLD, "!"),
        ("question", SOFT_GREY, SLATE, "× cache miss"),
    ]
    yy = 2.3
    for txt, fill, stroke, mark in stack:
        filled_rect(s, 7.4, yy, 4.0, 0.42, fill, stroke=stroke,
                    stroke_pt=1.4, radius=True, radius_adj=0.2)
        text_box(s, 7.55, yy + 0.02, 3.75, 0.38, txt,
                 size=10.5 if len(txt) > 20 else 12, bold=(mark == "!"),
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 11.5, yy + 0.02, 1.2, 0.38, mark, size=11, bold=True,
                 color=stroke if mark != "!" else GOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.52
    text_box(s, 7.4, 4.95, 5.15, 0.6,
             "A dynamic element at the start (timestamp) breaks the "
             "cache for every block after it — this is “request 3.”",
             size=11.5, italic=True, color=SLATE, line_spacing=1.12)
    gold_callout(s, 0.55, 5.95, 12.25, 0.8,
                 "Composition rule: stable content goes first (prompt, "
                 "instructions, examples, documents), variable content "
                 "goes last.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    """Window race: log-bars for 4 models + 2 outliers + a line about
    positional encoding + gold callout."""
    s = blank(p)
    slide_title(s, "The 2026 frontier standard is a window of up to "
                   "1 million tokens. But advertised ≠ usable",
                size=24, h=0.9, y=0.35)
    # Log bars
    ocean_box(s, 0.55, 1.5, 7.65, 3.55, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    bars = [
        ("GPT-3.5 (2022)", "4K", 4_000, MID, ""),
        ("Claude 3.5 (2024)", "200K", 200_000, MID, ""),
        ("Frontier standard (2026)", "up to 1M", 1_000_000, GOLD,
         "Fable 5, GPT-5.6, Gemini 3.1 Pro and others"),
        ("Gemini 3.5 Pro (2026)", "2M", 2_000_000, MID,
         "a single outlier model, not the standard"),
    ]
    import math
    yy = 1.68
    for name, val, v, color, note in bars:
        w = (math.log10(v) - 3.0) / 3.4 * 4.5
        text_runs(s, 0.9, yy, 6.9, 0.3, [
            {"text": name + " — ", "size": 12.5, "bold": True,
             "color": DEEP},
            {"text": val, "size": 12.5, "bold": True,
             "color": GOLD if color == GOLD else MID},
            {"text": ("  ·  " + note) if note else "", "size": 11,
             "italic": True, "color": SLATE},
        ])
        filled_rect(s, 0.9, yy + 0.33, max(w, 0.3), 0.26, color,
                    radius=True, radius_adj=0.35)
        yy += 0.78
    text_box(s, 0.9, 4.73, 6.9, 0.3, "width — logarithmic scale",
             size=10.5, italic=True, color=SLATE)
    # Outliers
    card = filled_rect(s, 8.45, 1.5, 4.35, 1.68, SOFT_GREY, stroke=LIGHT,
                       stroke_pt=1.4, radius=True, radius_adj=0.08)
    card.line.dash_style = 4
    text_box(s, 8.7, 1.65, 3.85, 0.4, "Llama 4 Scout: “10M”", size=14,
             bold=True, color=DEEP)
    text_box(s, 8.7, 2.1, 3.85, 1.0,
             "advertised; no published benchmark confirms quality near "
             "the limit", size=12, color=SLATE,
             line_spacing=1.18)
    ocean_box(s, 8.45, 3.37, 4.35, 1.68, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 8.7, 3.52, 3.85, 0.4, "YandexGPT 5 Pro: 32K", size=14,
             bold=True, color=DEEP)
    text_box(s, 8.7, 3.97, 3.85, 1.0,
             "one to two orders of magnitude below the flagships — for "
             "long documents this is a defining constraint", size=12,
             color=DEEP, line_spacing=1.18)
    text_box(s, 0.55, 5.2, 12.25, 0.55,
             "You can't just “stretch” the window: token position is "
             "encoded by a geometry trained on specific lengths — "
             "extending it (RoPE / YaRN) is separate engineering work.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER,
             line_spacing=1.15)
    gold_callout(s, 0.55, 5.9, 12.25, 0.68,
                 "You pay for what you put in the window, not for what "
                 "the window can hold: 900K input tokens at $10/million "
                 "≈ $9 for a single call.", size=13.5,
                 align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.63, 12.25, 0.55, [
        {"text": "What to do: ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": "choose a model by the task's effective window "
                 "(benchmarks without lexical shortcuts), not by the "
                 "marketing maximum.", "size": 12.5, "color": DEEP},
    ], align=PP_ALIGN.CENTER, line_spacing=1.12)
    speaker_notes(s, load_notes("s23"))


def build_s25(p):
    """Two-tier layout (top/bottom): top — "forgetting is solved"
    (U-curve 2023 -> flat line 2026, own chart); bottom — "understanding
    by meaning — no" (NoLiMa panel with explicit caption). + formula
    callout."""
    s = blank(p)
    slide_title(s, "Verbatim retrieval is nearly solved. Understanding "
                   "long context is not", size=23, h=0.5,
                y=0.30)
    # -- TOP TIER: forgetting is solved
    ocean_box(s, 0.55, 1.02, 12.25, 2.42, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.85, 1.12, 11.7, 0.4, [
        {"text": "Forgetting is solved ", "size": 14.5, "bold": True,
         "color": MID},
        {"text": "✓", "size": 16, "bold": True, "color": TEAL},
        {"text": "  — finding a verbatim insertion (needle-in-a-haystack)",
         "size": 13, "color": DEEP},
    ])
    add_image(s, ASSETS / "charts/s25-ucurve-en.png", x=0.85, y=1.56, w=6.2)
    text_runs(s, 7.3, 1.62, 2.72, 1.75, [
        {"text": "The 2023 U-curve has flattened: ", "size": 12,
         "color": DEEP},
        {"text": "single-needle — up to 99% on the full 1M window",
         "size": 12, "bold": True, "color": DEEP},
        {"text": " (Gemini Deep Think).", "size": 12, "color": DEEP},
    ], line_spacing=1.18)
    # Real photo of a needle in a haystack — a literal anchor metaphor
    # for needle-in-a-haystack, full frame without aggressive cropping.
    ocean_box(s, 10.08, 1.52, 2.68, 1.86, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.2)
    _place_image_contain(s, ASSETS / "web/needle-haystack-crop.jpg",
                         10.15, 1.58, 2.54, 1.74)
    # -- BOTTOM TIER: understanding by meaning — no
    ocean_box(s, 0.55, 3.58, 12.25, 2.85, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.85, 3.68, 11.7, 0.4, [
        {"text": "Understanding by meaning — no ", "size": 14.5,
         "bold": True, "color": DEEP},
        {"text": "(NoLiMa, 2025: the benchmark removed literal word "
                 "overlap)", "size": 12.5, "color": SLATE},
    ])
    add_image(s, ASSETS / "charts/s25-nolima-en.png", x=0.85, y=4.12, w=4.6)
    text_runs(s, 5.75, 4.25, 6.8, 1.7, [
        {"text": "Without lexical overlap, quality falls ", "size": 13,
         "color": DEEP},
        {"text": "below 50% of the baseline already at 32K", "size": 13.5,
         "bold": True, "color": GOLD},
        {"text": " — for 11 of 13 models.", "size": 13, "bold": True,
         "color": DEEP},
        {"text": "32K tokens ≈ 3% of a flagship's advertised window; "
                 "the baseline is the same model's accuracy on short "
                 "context.",
         "size": 11.5, "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 8},
    ], line_spacing=1.2)
    gold_callout(s, 0.55, 6.55, 12.25, 0.72,
                 "A 1M window ≠ 1M tokens of reasoning. The window is "
                 "how much the model can read; the usable length is how "
                 "many tokens out it can still connect facts across.",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s25"))


# ============================================================
# BATCH 2 — Section 4. Sampling and Generation
# ============================================================
def build_s26a(p):
    section_divider(
        p, section_n=4, sub_title="Sampling and Generation",
        frame_phrase="How a probability distribution gives birth to a "
                     "single token — and which knobs steer that choice",
        tag="4 case studies · 2 failures", active_stage={4, 5},
        notes_id="s26a",
        passed={0, 1, 2, 3},
        frame_size=18,
        illus=ASSETS / "web/dice-wikimedia.jpg")


S26_BARS = [("apple", 0.32, True), ("pizza", 0.19, False),
            ("salad", 0.14, False), ("a sandwich", 0.11, False),
            ("a cucumber", 0.08, False)]


def build_s26(p):
    """Top-5 distribution (native bars) + sampling -> one token +
    footnote + gold callout."""
    s = blank(p)
    slide_title(s, "At every step — a probability distribution over "
                   "the entire token vocabulary", size=24, h=0.9, y=0.35)
    # Left box — bars
    ocean_box(s, 0.55, 1.5, 7.9, 3.7, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.9, 1.7, 7.2, 0.45, [
        {"text": "Context: ", "size": 15, "bold": True, "color": MID},
        {"text": "“Today I ate …”", "size": 15.5, "bold": True,
         "color": DEEP, "font": FONT_MONO},
    ])
    text_box(s, 0.9, 2.2, 7.2, 0.35, "P(next token):", size=13,
             color=SLATE, italic=True)
    yy = 2.62
    for lab, v, is_gold in S26_BARS:
        text_box(s, 0.9, yy, 1.35, 0.4, lab, size=13.5, bold=is_gold,
                 color=DEEP, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.RIGHT)
        bw = v / 0.35 * 4.6
        filled_rect(s, 2.4, yy + 0.05, bw, 0.32,
                    GOLD if is_gold else MID, radius=True, radius_adj=0.3)
        text_runs(s, 2.5 + bw, yy, 1.6, 0.4, [
            {"text": f"{v:.2f}", "size": 13.5,
             "bold": True, "color": GOLD if is_gold else DEEP},
            {"text": "  — the maximum" if is_gold else "", "size": 11,
             "italic": True, "color": SLATE},
        ], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.5
    # Right box — sampling
    ocean_box(s, 8.75, 1.5, 4.05, 3.7, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 9.0, 1.7, 3.55, 0.45, "Sampling → one token", size=15,
             bold=True, color=TEAL)
    text_box(s, 9.0, 2.25, 3.55, 0.7,
             "the rule for picking one token from the distribution — "
             "the only knob in your hands",
             size=12, color=DEEP, line_spacing=1.2)
    down_arrow(s, 10.6, 3.1, w=0.4, h=0.5, fill=TEAL)
    chip(s, 9.55, 3.75, 2.4, 0.55, "[ apple ]", fill=GOLD, color=DEEP,
         size=15, font=FONT_MONO)
    text_box(s, 9.0, 4.42, 3.55, 0.6, "one is chosen — the remaining "
             "candidates vanish from the answer", size=11, italic=True,
             color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.15)
    text_box(s, 0.55, 5.38, 12.25, 0.4,
             "The remaining ~200,000 vocabulary tokens — each below "
             "0.05. All probabilities sum to 1. The numbers are "
             "illustrative.",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 5.85, 12.25, 1.0,
                 "The distribution is the model's “real” output. A "
                 "confident answer and a hallucination existed "
                 "simultaneously before sampling — as probability mass; "
                 "the policy made the choice.",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s26"))


def _mini_bars(s, x, y, w, h, vals, *, gold_first=False):
    """Mini distribution: vertical bars in area (x,y,w,h)."""
    n = len(vals)
    bw = w / (n * 1.5)
    gap = bw * 0.5
    vmax = max(vals) or 1.0
    xx = x + (w - (n * bw + (n - 1) * gap)) / 2
    for i, v in enumerate(vals):
        bh = max(v / vmax * h, 0.02)
        filled_rect(s, xx, y + h - bh, bw, bh,
                    GOLD if (gold_first and i == 0) else MID,
                    radius=False)
        xx += bw + gap
    plain_line(s, x, y + h + 0.02, x + w, y + h + 0.02, color=LIGHT,
               w_pt=1.2)


def build_s27(p):
    """Explicit formula (argmax / P^(1/T)) + 3 distribution panels +
    "why T matters" callout + top-p/top-k + live badge."""
    s = blank(p)
    slide_title(s, "Temperature is a divisor on the logits: sharpness "
                   "of choice, not knowledge",
                size=23, h=0.5, y=0.32)
    # Formula — explicit, in an Ocean rounded box
    ocean_box(s, 0.55, 0.98, 12.25, 0.85, fill=SURFACE, stroke=MID,
              stroke_pt=1.4)
    text_runs(s, 0.85, 1.04, 11.7, 0.75, [
        {"text": "T = 0  ⇒  choice = argmax P(token)", "size": 14.5,
         "bold": True, "color": DEEP, "font": FONT_MONO},
        {"text": "      T > 0  ⇒  sampling from a P^(1/T)-shaped "
                 "distribution", "size": 14.5, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": "logits are divided by T before softmax: T<1 sharpens "
                 "the distribution, T>1 flattens it", "size": 11.5,
         "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 4},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line_spacing=1.1)
    panels = [
        ("T → 0", "(argmax)", [1.0, 0.001, 0.001, 0.001, 0.001],
         "Picks the most likely token. Nearly identical answers — "
         "we'll unpack that “nearly” on the next slide.", False),
        ("T = 1", "(default)", [0.32, 0.19, 0.14, 0.11, 0.08],
         "Sampling proportional to the model's probabilities.", True),
        ("T = 1.5", "(smoothing)", [0.24, 0.19, 0.16, 0.14, 0.12],
         "Rare tokens get a real chance — anywhere up to incoherence.",
         False),
    ]
    x = 0.55
    for tname, tsub, vals, caption, is_std in panels:
        if is_std:
            filled_rect(s, x, 2.0, 3.95, 2.95, WHITE, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.05)
        else:
            ocean_box(s, x, 2.0, 3.95, 2.95, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.3)
        text_runs(s, x + 0.25, 2.14, 3.45, 0.42, [
            {"text": tname, "size": 16, "bold": True, "color": DEEP,
             "font": FONT_MONO},
            {"text": "  " + tsub, "size": 13.5, "bold": True,
             "color": GOLD if is_std else MID},
        ])
        _mini_bars(s, x + 0.55, 2.62, 2.85, 1.15, vals, gold_first=True)
        text_box(s, x + 0.25, 4.02, 3.45, 0.85, caption, size=11,
                 color=DEEP, line_spacing=1.15)
        x += 4.15
    # Why T matters even though token order doesn't change
    gold_callout(s, 0.55, 5.10, 12.25, 0.78,
                 "Why T matters even though the ranking of candidates "
                 "doesn't change: the choice isn't “take the top-1” but "
                 "a random draw proportional to probabilities; T "
                 "reshapes the probabilities themselves.", size=13,
                 align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.00, 12.3, 0.42, [
        {"text": "top-p", "size": 13, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": " — cuts the tail by probability mass · ", "size": 12.5,
         "color": DEEP},
        {"text": "top-k", "size": 13, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": " — by number of candidates · T per task: 0–0.3 code "
                 "and classification, 0.7+ generation.", "size": 12.5,
         "color": DEEP},
    ], align=PP_ALIGN.CENTER)
    ocean_box(s, 2.65, 6.52, 8.0, 0.55, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_runs(s, 2.85, 6.52, 7.6, 0.55, [
        {"text": "Live run: ", "size": 12.5, "bold": True,
         "color": TEAL},
        {"text": "one prompt — 10 times at T=0 and 10 times at T=1.5.",
         "size": 12.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """T=0 != determinism: 80/1000 exhibit + cause chain + 2 tiles +
    gold callout."""
    s = blank(p)
    slide_title(s, "T=0 does not give determinism: 80 unique answers "
                   "out of 1000", size=25, h=0.6)
    # Left — exhibit (height per content, frees space for the meme below)
    ocean_box(s, 0.55, 1.35, 4.55, 2.15, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.8, 1.48, 4.05, 0.85, "80 / 1000", size=48, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, 0.85, 2.28, 3.95, 1.15,
             "unique response variants to an identical request at T=0 "
             "— stock vLLM (an open inference server; Thinking Machines "
             "Lab, September 2025)",
             size=11, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.12)
    # Real "Well yes, but actually no" meme — same template as s01,
    # different text: T=0 sounds like a determinism guarantee, but no.
    meme_h = 1.08
    meme_w = meme_h * 1600 / 1218
    meme_x = 0.55 + (4.55 - meme_w) / 2
    meme_y = 3.60
    add_image(s, ASSETS / "web/well-yes-actually-no-template.jpg",
              x=meme_x, y=meme_y, h=meme_h)
    band_h = meme_h * 0.205
    text_box(s, meme_x + 0.08, meme_y + 0.02, meme_w - 0.16, band_h - 0.03,
             "T=0 = determinism?", size=10, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=0.95)
    # Right — mechanism
    text_runs(s, 5.4, 1.35, 7.4, 0.6, [
        {"text": "The cause isn't “floating-point in general” but a "
                 "lack of ", "size": 13.5, "bold": True, "color": DEEP},
        {"text": "batch invariance in the kernels:", "size": 13.5,
         "bold": True, "color": MID},
    ], line_spacing=1.15)
    chain = [
        "The server groups concurrent requests from different users "
        "into batches",
        "Batch size depends on other users' load in that exact "
        "millisecond",
        "Different batch size → different summation order → different "
        "least-significant bits",
        "Two close argmax candidates → the least-significant bit "
        "decides the token → autoregression spreads the divergence "
        "through the whole answer",
    ]
    yy = 2.0
    for i, step in enumerate(chain):
        h = 0.55 if i < 3 else 0.7
        ocean_box(s, 5.4, yy, 7.4, h, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, 5.6, yy, 7.0, h, [
            {"text": f"{i+1}. ", "size": 12, "bold": True, "color": TEAL},
            {"text": step, "size": 11.5, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        yy += h + 0.08
    # Two tiles
    filled_rect(s, 0.55, 4.75, 6.0, 1.0, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.3, radius=True, radius_adj=0.1)
    text_runs(s, 0.75, 4.82, 5.6, 0.9, [
        {"text": "A fix exists: ", "size": 12.5, "bold": True,
         "color": TEAL},
        {"text": "batch-invariant kernels — 1000/1000 bit-for-bit "
                 "identical", "size": 12.5, "color": DEEP},
    ], line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)
    ocean_box(s, 6.8, 4.75, 6.0, 1.0, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.2)
    text_runs(s, 7.0, 4.82, 5.6, 0.9, [
        {"text": "Cost: ~35% of throughput", "size": 12.5,
         "bold": True, "color": DEEP},
        {"text": " → providers don't turn it on; OpenAI's seed is "
                 "officially “mostly deterministic,” not fully",
         "size": 12, "color": DEEP},
    ], line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 6.05, 12.25, 0.95,
                 "You cannot get a guaranteed-deterministic answer from "
                 "a cloud LLM today — build your processes accounting "
                 "for that. Consequence: don't compare answers "
                 "bit-for-bit, compare them semantically or "
                 "structurally.",
                 size=14, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """Table "parameter — range — what it affects — typical value" (5
    knobs, effort — gold); budget_tokens/400-error kept speaker-notes
    only."""
    s = blank(p)
    slide_title(s, "API knobs: reasoning depth has joined randomness "
                   "and length", size=24, h=0.6)
    # 6x4 table
    tx, ty, tw = 0.55, 1.45, 12.25
    col_ws = [2.9, 1.85, 4.1, 3.4]
    gtbl = s.shapes.add_table(6, 4, Inches(tx), Inches(ty), Inches(tw),
                              Inches(4.35))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, w in enumerate(col_ws):
        tbl.columns[ci].width = Inches(w)
    tbl.rows[0].height = Inches(0.45)
    for ri in range(1, 6):
        tbl.rows[ri].height = Inches(0.78)

    def cell(r, c, txt, *, size=12, bold=False, color=DEEP, fill=WHITE,
             mono=False, align=PP_ALIGN.LEFT):
        cl = tbl.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
        cl.margin_left = Inches(0.08); cl.margin_right = Inches(0.06)
        cl.margin_top = Inches(0.02); cl.margin_bottom = Inches(0.02)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cl.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        para.line_spacing = 1.08
        r_ = para.add_run(); r_.text = txt
        r_.font.name = FONT_MONO if mono else FONT_BODY
        r_.font.size = Pt(size); r_.font.bold = bold
        r_.font.color.rgb = color

    headers = ["Parameter", "Range", "What it affects",
               "Typical value"]
    for ci, htxt in enumerate(headers):
        cell(0, ci, htxt, size=12.5, bold=True, color=MID)
    rows_data = [
        ("temperature", "0–2", "Determinism ↔ chaos of token choice",
         "0 classification; 0.7–1.2 text", False),
        ("top_p", "0.1–1", "Width of the candidate tail during sampling",
         "0.9–0.95", False),
        ("max_tokens", "integer", "Hard cutoff on generation — can stop "
         "mid-way", "set per task, with headroom for JSON/code", False),
        ("effort", "none → xhigh", "Depth of internal reasoning — "
         "and its cost", "medium by default for most providers", True),
        ("verbosity", "low → high", "Length of the visible answer, "
         "independent of thinking depth", "medium", False),
    ]
    for ri, (name, rng, eff, typ, is_gold) in enumerate(rows_data):
        fill = GOLD_TINT if is_gold else (SURFACE if ri % 2 == 0 else WHITE)
        cell(ri + 1, 0, name, size=13, bold=True,
             color=DEEP if is_gold else MID, fill=fill, mono=True)
        cell(ri + 1, 1, rng, size=12, fill=fill, mono=True,
             align=PP_ALIGN.CENTER)
        cell(ri + 1, 2, eff, size=12, fill=fill, bold=is_gold)
        cell(ri + 1, 3, typ, size=11.5, fill=fill)
    text_box(s, 0.55, 5.95, 12.25, 0.35,
             "effort / reasoning_effort — a new axis for 2026: OpenAI — "
             "an effort scale; Anthropic — adaptive thinking; Gemini — "
             "thinking budget.", size=11.5, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.40, 12.25, 0.68,
                 "What to do: start tuning with temperature and effort "
                 "— these are the two main knobs; top_p / top_k / "
                 "verbosity are fine-tuning on top.", size=13.5,
                 align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """Constrained decoding: masking mechanics + limits + retrieval
    pause."""
    s = blank(p)
    slide_title(s, "Structured outputs: invalid tokens are zeroed out "
                   "right in the distribution", size=23, h=0.9, y=0.35)
    # Left — mechanics
    ocean_box(s, 0.55, 1.5, 6.9, 4.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.85, 1.65, 6.3, 1.15, [
        {"text": "Constrained decoding: ", "size": 13, "bold": True,
         "color": MID},
        {"text": "the JSON schema is compiled into a ", "size": 13,
         "color": DEEP},
        {"text": "finite-state machine over tokens", "size": 13,
         "bold": True, "color": DEEP},
        {"text": ". During token-by-token generation, the automaton "
                 "tracks the state of the prefix so far and, at every "
                 "step, MASKS (zeroes the probability of) tokens that "
                 "would lead to an invalid continuation.",
         "size": 13, "color": DEEP},
    ], line_spacing=1.18)
    # Mini masking diagram: 5 bars, 2 disabled
    vals = [(0.32, True), (0.19, False), (0.14, True), (0.11, False),
            (0.08, True)]
    xx = 1.45
    for v, valid in vals:
        bh = v / 0.35 * 1.0
        filled_rect(s, xx, 2.95 + (1.0 - bh), 0.6, bh,
                    MID if valid else SOFT_GREY, radius=False)
        if not valid:
            text_box(s, xx + 0.08, 2.82 + (1.0 - bh) - 0.28, 0.5, 0.35, "×",
                     size=17, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        xx += 1.0
    plain_line(s, 1.25, 3.98, 6.45, 3.98, color=LIGHT, w_pt=1.2)
    text_box(s, 1.25, 4.04, 5.4, 0.3,
             "the automaton masks the gray bars → probability 0",
             size=11, italic=True, color=SLATE)
    text_runs(s, 0.85, 4.48, 6.3, 1.1, [
        {"text": "Asking “answer strictly in JSON” → ~80% valid\n",
         "size": 13, "color": DEEP},
        {"text": "Strict mode → ", "size": 13.5, "bold": True,
         "color": DEEP, "newpara": True, "space_before_pt": 4},
        {"text": "100%", "size": 15, "bold": True, "color": GOLD},
        {"text": " — the output is valid by construction, not "
                 "“checked after the fact”", "size": 12.5, "color": DEEP},
    ], line_spacing=1.18)
    # Right — limitations
    ocean_box(s, 7.75, 1.5, 5.05, 4.25, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 8.0, 1.65, 4.55, 0.4,
             "Limitations — properties of compilation", size=14,
             bold=True, color=MID)
    limits = [
        "Recursion via $ref — not supported (tree → flat list with "
        "parent_id)",
        "Nesting depth — ≤ 5",
        "The first request with a new schema pays for grammar "
        "compilation — up to 10 s",
        "Guarantees syntax, not meaning: validating values is still "
        "your job",
    ]
    yy = 2.15
    for lim in limits:
        filled_rect(s, 8.0, yy, 4.55, 0.76, SURFACE, stroke=LIGHT,
                    stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(s, 8.15, yy + 0.04, 4.25, 0.68, lim, size=11, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        yy += 0.87
    # "Ask the room" moved left, right flank — Gandalf "INVALID TOKEN
    # SHALL NOT PASS" (direct masking metaphor).
    ocean_box(s, 0.55, 6.15, 8.6, 0.65, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_runs(s, 0.75, 6.15, 8.2, 0.65, [
        {"text": "Ask the room: ", "size": 14, "bold": True, "color": TEAL},
        {"text": "why exactly 100%, and not 99.9?", "size": 14,
         "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    _place_image_contain(s, ASSETS / "web/gandalf-token-en.jpg",
                         9.4, 5.92, 3.2, 1.40)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """Reframed around stop conditions: 5-step loop + return arrow +
    cards "how it stops" / "how it breaks" (degenerate repetition, "Pot,
    stop boiling!" meme) + practice (repetition penalty, max_tokens)."""
    s = blank(p)
    slide_title(s, "Predict a token → append it to context → predict "
                   "the next one", size=21, h=0.55)
    steps = [
        ("1 · Current context",
         "system prompt + history + request + the already-generated "
         "part of the answer", False),
        ("2 · Forward pass",
         "tokenization → embeddings → all attention layers", True),
        ("3 · Distribution",
         "probabilities over all ~200K tokens", False),
        ("4 · Sampling",
         "one token — temperature / top-p / schema", False),
        ("5 · Token appended",
         "to context — and the loop repeats", False),
    ]
    bw, gap = 2.32, 0.24
    x0 = (SLIDE_W_IN - bw * 5 - gap * 4) / 2
    y0, bh = 1.72, 1.55
    for i, (head, desc, is_gold) in enumerate(steps):
        x = x0 + i * (bw + gap)
        if is_gold:
            filled_rect(s, x, y0, bw, bh, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.08)
        else:
            ocean_box(s, x, y0, bw, bh, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.3)
        text_box(s, x + 0.15, y0 + 0.10, bw - 0.3, 0.62, head, size=12.5,
                 bold=True, color=DEEP, line_spacing=1.08)
        text_box(s, x + 0.15, y0 + 0.72, bw - 0.3, 0.78, desc, size=10.5,
                 color=MID if not is_gold else DEEP, line_spacing=1.12)
        if i < 4:
            right_arrow(s, x + bw + 0.02, y0 + bh / 2 - 0.09, w=gap - 0.04,
                        h=0.18, fill=MID)
    chip(s, x0 - 0.15, y0 - 0.40, 0.9, 0.32, "input", fill=GOLD, color=DEEP,
         size=11.5)
    # Return arrow
    plain_line(s, x0 + 4 * (bw + gap) + bw / 2, y0 + bh,
               x0 + 4 * (bw + gap) + bw / 2, y0 + bh + 0.30, color=LIGHT,
               w_pt=2.0)
    left_arrow(s, x0 + bw / 2, y0 + bh + 0.24, w=4 * (bw + gap), h=0.18,
               fill=LIGHT)
    plain_line(s, x0 + bw / 2, y0 + bh, x0 + bw / 2, y0 + bh + 0.26,
               color=LIGHT, w_pt=2.0)
    text_box(s, x0 + 3.0, y0 + bh + 0.44, 6.0, 0.32,
             "⟲ return to step 1 — the loop repeats", size=11.5,
             italic=True, color=MID, align=PP_ALIGN.CENTER)
    # How the loop stops / how it breaks
    ocean_box(s, 0.55, 4.15, 5.95, 1.75, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_box(s, 0.85, 4.27, 5.4, 0.38, "How the loop stops",
             size=14, bold=True, color=TEAL)
    text_runs(s, 0.85, 4.68, 5.4, 1.1, [
        {"text": "A special ", "size": 12.5, "color": DEEP},
        {"text": "stop token", "size": 12.5, "bold": True, "color": DEEP},
        {"text": " — the model itself decides “answer finished” — or ",
         "size": 12.5, "color": DEEP},
        {"text": "max_tokens", "size": 12.5, "bold": True, "color": DEEP,
         "font": FONT_MONO},
        {"text": ": the cutoff is instant, even mid-way through a JSON "
                 "field.", "size": 12.5, "color": DEEP},
    ], line_spacing=1.18)
    ocean_box(s, 6.8, 4.15, 6.0, 1.75, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 7.1, 4.27, 4.0, 0.38, "How the loop breaks",
             size=14, bold=True, color=DEEP)
    text_box(s, 7.1, 4.68, 3.05, 1.1,
             "A degenerate repetition loop: the model “gets stuck” on "
             "one token or phrase and generates a wall of repeats "
             "instead of stopping.", size=12, color=DEEP,
             line_spacing=1.15)
    # Frame from the Soviet cartoon "Pot of Porridge" (Soyuzmultfilm,
    # 1984) — the "pot, stop boiling!" moment, a direct metaphor for a
    # loop with no stop condition. Pure crop, no baked text.
    gsh_h = 1.50
    gsh_w = gsh_h * 1040 / 720
    gsh_x = 12.55 - gsh_w
    gsh_y = 4.30 + (1.5 - gsh_h) / 2 + 0.08
    add_image(s, ASSETS / "web/gorshochek-1984-crop.jpg",
              x=gsh_x, y=gsh_y, h=gsh_h)
    gold_callout(s, 0.55, 6.05, 12.25, 0.68,
                 "Practice: repetition_penalty / frequency_penalty "
                 "reduce the probability of literal repeats; "
                 "max_tokens is a safety net against an infinite loop.",
                 size=13, align=PP_ALIGN.CENTER)
    text_box(s, 0.55, 6.82, 12.3, 0.4,
             "Every step is stateless: all “memory” lives in the "
             "context, which is fed in whole each time (the KV-cache "
             "makes re-feeding it cheap, without changing this "
             "logically).",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    """Reasoning tokens: forked output + pricing bars + boundaries +
    gold callout."""
    s = blank(p)
    slide_title(s, "Reasoning tokens are invisible but billed as output",
                size=25, h=0.6)
    # Top — forked output
    ocean_box(s, 0.55, 1.3, 12.25, 2.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    filled_rect(s, 0.95, 1.9, 2.7, 0.75, MID, radius=True, radius_adj=0.12)
    text_box(s, 1.0, 1.95, 2.6, 0.65, "autoregressive\nloop  ⟲",
             size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    line_arrow(s, 3.7, 2.1, 4.4, 1.8, color=TEAL, w_pt=2.2)
    line_arrow(s, 3.7, 2.45, 4.4, 2.8, color=SLATE, w_pt=2.2)
    filled_rect(s, 4.5, 1.55, 4.6, 0.5, TEAL, radius=True, radius_adj=0.2)
    text_box(s, 4.65, 1.58, 4.3, 0.44, "visible answer", size=12.5,
             bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 4.5, 2.55, 7.7, 0.62, SOFT_GREY, radius=True,
                radius_adj=0.14)
    text_box(s, 4.7, 2.55, 6.2, 0.62,
             "a draft “for itself” — doesn't reach the answer, does "
             "reach the bill",
             size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.1)
    text_box(s, 11.05, 2.5, 1.15, 0.72, "×3–10", size=19, bold=True,
             color=GOLD, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # Joker burning a pile of money (The Dark Knight) — direct
    # reinforcement of "invisible tokens burn budget," right above
    # "x3-10".
    _place_image_contain(s, ASSETS / "web/joker-burning-money.jpg",
                         10.42, 1.42, 2.3, 1.07)
    text_box(s, 4.5, 3.24, 8.2, 0.3,
             "at the output-token rate, with a natural ceiling in "
             "max_tokens — but no built-in ceiling of its own",
             size=11, italic=True, color=SLATE)
    # Bottom left — pricing bars
    ocean_box(s, 0.55, 3.75, 5.95, 2.2, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    pb = [("o4-mini", 1, MID), ("o3", 5, MID), ("o3-pro", 18, GOLD)]
    yy = 3.95
    for lab, v, color in pb:
        text_box(s, 0.8, yy, 1.05, 0.32, lab, size=11.5, bold=True,
                 color=DEEP, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        bw = v / 18 * 3.6
        filled_rect(s, 1.95, yy + 0.04, max(bw, 0.18), 0.24, color,
                    radius=True, radius_adj=0.35)
        text_box(s, 2.05 + max(bw, 0.18), yy, 0.9, 0.32, f"{v}×",
                 size=12, bold=True, color=GOLD if color == GOLD else DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.42
    text_runs(s, 0.8, 5.28, 5.5, 0.6, [
        {"text": "o3-pro: 3.6× more expensive than o3, 18× more "
                 "expensive than o4-mini", "size": 11.5,
         "bold": True, "color": DEEP},
        {"text": " — at comparable visible-answer length; the "
                 "difference comes from the volume of reasoning.",
         "size": 11, "color": SLATE},
    ], line_spacing=1.12)
    # Bottom right — 2 boundaries
    borders = [
        ("Control:", " adaptive thinking / effort instead of manual "
         "budgets — convenient, but request cost has become less "
         "predictable"),
        ("The “chain of thought” in the UI is a paraphrase",
         " (summarized), not raw tokens: you cannot build a decision "
         "audit on “displayed thoughts”"),
    ]
    yy = 3.75
    for head, rest in borders:
        ocean_box(s, 6.75, yy, 6.05, 1.02, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, 6.95, yy + 0.06, 5.65, 0.9, [
            {"text": head, "size": 12, "bold": True, "color": MID},
            {"text": rest, "size": 11.5, "color": DEEP},
        ], line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.18
    gold_callout(s, 0.55, 6.15, 12.25, 0.85,
                 "Budget for an invisible portion 2–5× the visible "
                 "answer — and cross-check against the usage field in "
                 "the API response, where reasoning tokens show up as "
                 "a line item.",
                 size=14, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s32"))


# ============================================================
# Section 5. Model Types and Sizes (NEW)
# ============================================================
def build_s33a(p):
    """Divider for the new Section 5. The whole pipeline strip is in
    muted gold (the section sits logically "above" the pipeline —
    classifying the models that run it; NOT a separate stage). No gold
    frame (the frame marks a full pass on s35a)."""
    section_divider(
        p, section_n=5, sub_title="Model Types and Sizes",
        frame_phrase="What models run on, of different sizes — and "
                     "what each class can do",
        tag="4 case studies", active_stage=set(), notes_id="s33a",
        bar_muted=True, frame_size=18,
        illus=ASSETS / "web/matryoshka-wikimedia.jpg")


S33_COLS = [
    ("s33-laptop", "Small — up to 8–10B", False, [
        ("Examples: ", "Qwen3.8-4B / 8B, Llama-class models"),
        ("Hardware: ", "laptop, smartphone, edge device"),
        ("Multimodality: ", "usually text-only or basic vision"),
    ]),
    ("s33-gpu", "Medium — around 30B", False, [
        ("Examples: ", "Muse Glimmer 30B — top of the medium class"),
        ("Hardware: ", "a single 24–32 GB GPU"),
        ("Multimodality: ", "often has vision"),
    ]),
    ("s33-server", "Large — 70B+", False, [
        ("Examples: ", "Llama-class 70B models, Qwen3.5-397B-A17B"),
        ("Hardware: ", "multi-GPU / server"),
        ("Multimodality: ", "generally full (text + image, sometimes "
         "audio)"),
    ]),
    ("s33-cloud", "MoE giants — 400B+", True, [
        ("Examples: ", "DeepSeek V4-Pro (1.6 trillion), Kimi K3 "
         "(2.8 trillion)"),
        ("Hardware: ", "cloud or a dedicated cluster ONLY — they don't "
         "fit on consumer hardware"),
        ("Multimodality: ", "full, top quality"),
    ]),
]


def build_s33(p):
    """Classification of models by size — 4-column matrix (a hardware
    icon atop each column), MoE giants gold + callout about memory as
    the limiting factor."""
    s = blank(p)
    slide_title(s, "Four size classes — from laptop to cluster",
                size=25, h=0.6)
    col_w, gap = 3.02, 0.14
    x0 = (SLIDE_W_IN - col_w * 4 - gap * 3) / 2
    y0, col_h = 1.28, 4.15
    for i, (icon, title, is_gold, rows) in enumerate(S33_COLS):
        x = x0 + i * (col_w + gap)
        if is_gold:
            filled_rect(s, x, y0, col_w, col_h, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.05)
        else:
            ocean_box(s, x, y0, col_w, col_h, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.4)
        add_image(s, ASSETS / f"icons/{icon}.png",
                  x=x + col_w / 2 - 0.27, y=y0 + 0.16, w=0.54)
        text_box(s, x + 0.1, y0 + 0.78, col_w - 0.2, 0.62, title,
                 size=14, bold=True, color=DEEP if is_gold else MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.05)
        runs = []
        for j, (head, rest) in enumerate(rows):
            runs.append({"text": head, "size": 12.5, "bold": True,
                         "color": TEAL if not is_gold else MID,
                         "newpara": j > 0, "space_before_pt": 14})
            runs.append({"text": rest, "size": 12.5, "color": DEEP,
                         "bold": is_gold and head.startswith("Hardware")})
        text_runs(s, x + 0.2, y0 + 1.52, col_w - 0.4, col_h - 1.65, runs,
                  line_spacing=1.22)
    gold_callout(s, 0.55, 5.50, 12.25, 0.95,
                 "The limiting factor is memory capacity, not compute. "
                 "The bigger the model, the wider its multimodality and "
                 "quality — and the smaller the odds of running it "
                 "yourself.",
                 size=14, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s33"))


def build_s34(p):
    """Local vs cloud 2026: three categories; middle one is the
    categorical news (gold)."""
    s = blank(p)
    slide_title(s, "“Open weights” no longer means “locally "
                   "runnable”", size=26, h=0.6)
    cols = [
        ("Truly local — up to ~30B", None, [
            "Qwen3.8-27B (Apache 2.0, image+video input, 262K window), "
            "Muse Glimmer 30B",
            "Hardware: RTX 5090 (32 GB) · Apple unified 64–128 GB",
            "The limit is memory capacity, not compute"]),
        ("Open-but-cloud-only giants", "open ≠ local", [
            "Kimi K3 — 2.8 trillion parameters, the largest open model",
            "DeepSeek V4-Pro — 1.6 trillion",
            "Don't fit on consumer hardware in any form"]),
        ("Closed APIs", None, [
            "Flagship-level quality — still lives here",
            "Pay per token, data flows through the provider"]),
    ]
    x = 0.55
    for title, badge, bullets in cols:
        is_mid = badge is not None
        if is_mid:
            filled_rect(s, x, 1.5, 3.95, 4.0, WHITE, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.05)
            chip(s, x + 0.55, 1.28, 2.85, 0.4, badge, fill=GOLD, color=DEEP,
                 size=13)
        else:
            ocean_box(s, x, 1.5, 3.95, 4.0, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.4)
        # First heading — one line, no wrap
        t_sz = 13.5 if "~30B" in title else 15
        text_box(s, x + 0.15, 1.82, 3.75, 0.75, title, size=t_sz, bold=True,
                 color=MID if not is_mid else DEEP, line_spacing=1.1)
        runs = []
        for i, b in enumerate(bullets):
            bold = b.startswith(("The limit", "Don't fit"))
            runs.append({"text": "• " + b, "size": 12.5, "color": DEEP,
                         "bold": bold, "newpara": i > 0,
                         "space_before_pt": 8})
        text_runs(s, x + 0.25, 2.7, 3.45, 2.6, runs, line_spacing=1.2)
        x += 4.15
    text_box(s, 0.55, 5.85, 12.25, 0.55,
             "The reasons for choosing local haven't changed: data "
             "privacy · no per-token cost at volume · independence "
             "from the network.",
             size=12.5, italic=True, color=MID, align=PP_ALIGN.CENTER,
             line_spacing=1.15)
    gold_callout(s, 1.4, 6.45, 10.55, 0.72,
                 "What to do: if the task fits in ~30B and the data "
                 "can't leave your perimeter — go local; if you need "
                 "flagship quality — go with a closed API.",
                 size=13, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s34"))


# ============================================================
# Section 6. Wrap-up
# ============================================================
def build_s35a(p):
    """Pepe Silvia / conspiracy board (It's Always Sunny) — a
    recognizable image of "pulling all the threads together," matching
    the theme "the pipeline is fully assembled"."""
    section_divider(
        p, section_n=6, sub_title="Wrap-up",
        frame_phrase="The pipeline assembled as a whole: a summary of "
                     "mechanisms and boundaries — and deciding when an "
                     "LLM is not the right tool",
        tag="6 case studies", active_stage=set(range(7)),
        notes_id="s35a", frame_bar=True, frame_size=18,
        illus=ASSETS / "web/pepe-silvia.jpg")


def build_s35(p):
    """Recap frame: "We've covered how the model works — let's assemble
    the picture"."""
    s = blank(p)
    slide_title(s, "We've covered how the model works — let's assemble "
                   "the picture", size=24, h=0.6)
    text_box(s, 0.55, 1.02, 12.3, 0.4,
             "New topics didn't add stages — they slotted into the "
             "existing ones.",
             size=14, italic=True, color=MID)
    stages = ["Tokenization", "Embeddings", "Attention", "Sampling"]
    bw, gap, y0, bh = 2.28, 0.3, 3.3, 0.85
    x0 = 0.62
    centers = []
    for i, st in enumerate(stages):
        x = x0 + i * (bw + gap)
        centers.append(x + bw / 2)
        ocean_box(s, x, y0, bw, bh, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.5)
        text_box(s, x, y0 + 0.08, bw, bh - 0.16, st, size=14.5, bold=True,
                 color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            right_arrow(s, x + bw + 0.03, y0 + bh / 2 - 0.09, w=gap - 0.06,
                        h=0.18, fill=MID)
    rx = x0 + 4 * (bw + gap)
    filled_rect(s, rx, y0 + 0.08, 1.85, bh - 0.16, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.8, radius=True, radius_adj=0.15)
    text_runs(s, rx, y0 + 0.12, 1.85, bh - 0.24, [
        {"text": "next", "size": 12, "bold": True, "color": DEEP,
         "align": PP_ALIGN.CENTER},
        {"text": "token", "size": 12, "bold": True, "color": DEEP,
         "newpara": True, "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    right_arrow(s, rx - gap + 0.03, y0 + bh / 2 - 0.09, w=gap - 0.06,
                h=0.18, fill=MID)
    # loop back
    plain_line(s, rx + 0.92, y0 + bh - 0.08, rx + 0.92, y0 + bh + 0.32,
               color=LIGHT, w_pt=1.8)
    left_arrow(s, centers[0], y0 + bh + 0.26, w=rx + 0.92 - centers[0],
               h=0.16, fill=LIGHT)
    text_box(s, 4.2, y0 + bh + 0.46, 5.0, 0.3, "⟲ loop — the token is "
             "appended to context", size=11, italic=True, color=MID,
             align=PP_ALIGN.CENTER)
    # Overlay tiles: 2 on top, 2 on the bottom
    overlays_top = [
        (0, 0.62, 3.35, "Glitch tokens", " — a vocabulary property at "
         "the tokenization stage"),
        (2, 4.9, 4.5, "KV-cache", " — inside attention; prompt caching "
         "— an add-on on top of it at the request boundary"),
    ]
    for stage_i, px, pw, head, rest in overlays_top:
        filled_rect(s, px, 1.55, pw, 0.95, TEAL_TINT, stroke=GOLD,
                    stroke_pt=1.5, radius=True, radius_adj=0.1)
        text_runs(s, px + 0.18, 1.60, pw - 0.36, 0.85, [
            {"text": head, "size": 12.5, "bold": True, "color": DEEP},
            {"text": rest, "size": 11.5, "color": DEEP},
        ], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
        plain_line(s, px + pw / 2, 2.5, centers[stage_i], y0, color=TEAL,
                   w_pt=1.4, dash=4)
    overlays_bottom = [
        (3, 7.6, 4.0, "Structured outputs", " — a filter at the "
         "sampling stage", centers[3]),
        (None, 2.9, 4.3, "Reasoning tokens", " — the same loop; part of "
         "the output marked as a “draft”", centers[0] + 1.0),
    ]
    for stage_i, px, pw, head, rest, tx_ in overlays_bottom:
        filled_rect(s, px, 5.0, pw, 0.95, TEAL_TINT, stroke=GOLD,
                    stroke_pt=1.5, radius=True, radius_adj=0.1)
        text_runs(s, px + 0.18, 5.05, pw - 0.36, 0.85, [
            {"text": head, "size": 12.5, "bold": True, "color": DEEP},
            {"text": rest, "size": 11.5, "color": DEEP},
        ], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
        plain_line(s, px + pw / 2, 5.0,
                   tx_ if stage_i is None else centers[stage_i],
                   y0 + bh + (0.3 if stage_i is None else 0),
                   color=TEAL, w_pt=1.4, dash=4)
    gold_callout(s, 0.55, 6.25, 12.25, 0.85,
                 "The pipeline is a diagnostic tree: by the symptom, "
                 "you can almost always guess the culprit stage. Ask: "
                 "“at which stage does this happen?”", size=14.5,
                 align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s35"))


def build_s36(p):
    """The 2026 landscape: 2 camps + IMO exhibit + price scale."""
    s = blank(p)
    slide_title(s, "September 2026: quality has converged — prices "
                   "have diverged by three orders of magnitude",
                size=24, h=0.9, y=0.35)
    # Columns — real company logos (LobeHub icons-static-svg, recolored
    # into the Ocean palette) as a visual anchor before each row.
    frontier = [
        ("openai", "OpenAI: GPT-5.6 — Luna → Terra → Sol"),
        ("anthropic", "Anthropic: Claude Fable 5 · Opus 5"),
        ("google", "Google: Gemini 3.5 Pro (2M window, Deep Think)"),
        ("xai", "xAI: Grok 4.3"),
    ]
    open_w = [
        ("deepseek", "DeepSeek V4 (Pro 1.6T / Flash 284B)"),
        ("qwen", "Qwen 3.8-Max — the first open model in the Max "
                 "lineup"),
        (None, "Kimi K2.6 (1T) · Kimi K3 (2.8T — the largest open "
               "model)"),
    ]
    for x, title, items in [(0.55, "Frontier (closed weights)",
                             frontier),
                            (6.8, "Open weights", open_w)]:
        ocean_box(s, x, 1.5, 6.0, 2.4, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.4)
        text_box(s, x + 0.25, 1.65, 5.5, 0.4, title, size=15, bold=True,
                 color=MID)
        row_y = 2.14
        row_h = 1.7 / len(items)
        for logo, it in items:
            bold = "largest open" in it
            if logo:
                add_image(s, ASSETS / f"icons/logos-web/{logo}.png",
                          x=x + 0.25, y=row_y + row_h / 2 - 0.14, h=0.28)
                tx0 = x + 0.62
            else:
                tx0 = x + 0.25
            text_box(s, tx0, row_y, x + 5.75 - tx0, row_h, it, size=12.5,
                     bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.1)
            row_y += row_h
    # IMO exhibit
    filled_rect(s, 0.55, 4.05, 12.25, 0.92, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.8, radius=True, radius_adj=0.12)
    text_box(s, 0.85, 4.13, 11.7, 0.4,
             "IMO 2026: six models — a perfect 42/42. Out of 666 human "
             "contestants — seven.", size=14.5, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, 0.85, 4.55, 11.7, 0.38,
             "The same systems miscount the letters in the word "
             "cranberry — “jagged intelligence” as a working "
             "characteristic.",
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)
    # Price scale
    filled_rect(s, 1.0, 5.55, 11.3, 0.2, SOFT_GREY, radius=True,
                radius_adj=0.5)
    filled_rect(s, 1.0, 5.55, 2.3, 0.2, LIGHT, radius=True, radius_adj=0.5)
    filled_rect(s, 10.0, 5.55, 2.3, 0.2, DEEP, radius=True, radius_adj=0.5)
    text_box(s, 1.0, 5.20, 4.5, 0.32, "market floor $0.03–0.2 / M tokens",
             size=12, bold=True, color=MID)
    text_box(s, 8.3, 5.20, 4.0, 0.32, "premium $10 in / $50 out",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.RIGHT)
    text_runs(s, 0.55, 5.98, 12.3, 0.4, [
        {"text": "Kimi K2.6 ≈ GPT-5.5 on the guarded SWE-bench Pro — ",
         "size": 13, "bold": True, "color": DEEP},
        {"text": "at ~80% less", "size": 13, "bold": True,
         "color": GOLD},
    ], align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.55, 12.25, 0.68,
                 "What to do: revisit your model choice regularly — "
                 "the landscape shifts on a scale of months, not "
                 "years.",
                 size=13, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s36"))


def build_s37(p):
    """Benchmarks: 3 story cards + gold callout."""
    s = blank(p)
    slide_title(s, "Benchmarks: contamination, overfitting — and "
                   "models that cheat on their own", size=24, h=0.9,
                y=0.35)
    # Story text <=2 lines (details in notes); 87.6 vs 57 as a large
    # contrasting stat block.
    cards = [
        ("1", "Contamination: memorized, not mastered", 8.2, [
            {"text": "SWE-bench: public repositories (Verified) vs "
                     "private codebases (Pro) — ", "size": 12.5,
             "color": DEEP},
            {"text": "the gap is the size of the memorization.",
             "size": 12.5, "bold": True, "color": DEEP},
            {"text": " OpenAI stopped publishing Verified in 2026.",
             "size": 12.5, "color": DEEP}]),
        ("2", "Overfitting to the leaderboard", 8.1, [
            {"text": "Llama 4 Maverick: on Chatbot Arena — a special "
                     "version, Elo 1417; the public model — ", "size": 12.5,
             "color": DEEP},
            {"text": "ranks 32–35", "size": 12.5, "bold": True,
             "color": DEEP},
            {"text": ". Yann LeCun: the results were “slightly cooked.”",
             "size": 12.5, "color": DEEP}]),
        ("3", "Models cheat on their own", 11.0, [
            {"text": "UK AI Security Institute: ", "size": 12.5,
             "color": DEEP},
            {"text": "all 5", "size": 12.5, "bold": True, "color": DEEP},
            {"text": " frontier models attempted to game the "
                     "evaluation process; one OpenAI model ", "size": 12.5,
             "color": DEEP},
            {"text": "escaped its sandbox and breached production "
                     "servers at Hugging Face", "size": 12.5,
             "bold": True, "color": DEEP},
            {"text": ".", "size": 12.5, "color": DEEP}]),
    ]
    yy = 1.52
    for num, head, body_w, runs in cards:
        ocean_box(s, 0.55, yy, 12.25, 1.42, fill=SURFACE if num != "3"
                  else WHITE, stroke=LIGHT, stroke_pt=1.3)
        text_box(s, 0.8, yy + 0.1, 0.55, 0.6, num, size=26, bold=True,
                 color=LIGHT)
        text_box(s, 1.45, yy + 0.12, 10.9, 0.38, head, size=14.5,
                 bold=True, color=MID)
        text_runs(s, 1.45, yy + 0.52, body_w, 0.85, runs, line_spacing=1.16)
        if num == "1":
            # Stat block: 87.6% (gold) vs 57% — large, contrasting
            text_runs(s, 9.85, yy + 0.16, 2.8, 0.75, [
                {"text": "87.6%", "size": 24, "bold": True, "color": GOLD},
                {"text": " vs ", "size": 14, "color": SLATE},
                {"text": "57%", "size": 24, "bold": True, "color": DEEP},
            ], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            text_box(s, 8.85, yy + 0.94, 3.8, 0.4,
                     "vendor claim vs guarded · average ~25%",
                     size=10.5, italic=True, color=SLATE,
                     align=PP_ALIGN.RIGHT)
        if num == "2":
            # Press X to Doubt (L.A. Noire) — reaction to showroom
            # leaderboard results, right third of the card.
            _place_image_contain(s, ASSETS / "web/pressx-template.jpg",
                                 10.15, yy + 0.07, 2.35, 1.28)
        yy += 1.55
    gold_callout(s, 0.55, 6.25, 12.25, 0.8,
                 "Benchmarks narrow the shortlist. Your own eval set "
                 "decides: 30–50 examples from your real tasks.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s37"))


# "Let's wrap up" — summary table "mechanism -> boundary -> what to do"
# across all topics actually covered (no M-codes, no checklist).
S38_ROWS = [
    ("Tokenization",
     "The model sees tokens, not letters; a viral fix doesn't "
     "generalize to the task class",
     "Test your own domain's “cranberry”; hand off counting to a tool"),
    ("Attention and role",
     "A role is tokens with weight; it shifts style and focus, not "
     "quality",
     "Role is a tone/focus tool; for facts, give data"),
    ("KV-cache and prompt caching",
     "Only saves on a repeated prefix; a token near the start breaks "
     "everything after it",
     "Stable content first, variable content last; watch cache hit"),
    ("Context window",
     "Advertised ≠ usable: 11 of 13 models lose half already at 32K",
     "Choose by benchmarks without lexical shortcuts"),
    ("Determinism at T=0",
     "Kernels aren't batch-invariant: someone else's load changes your "
     "answer",
     "Tests — not on bit-for-bit comparison; build processes "
     "accordingly"),
    ("Reasoning tokens",
     "Billed as output, with no natural ceiling",
     "Budget separately; set effort/verbosity explicitly"),
    ("Structured output",
     "Valid by construction — but not substantive quality",
     "Don't over-constrain the schema; you still validate the values"),
    ("Benchmarks",
     "Contamination, overfitting of showroom versions, models cheating",
     "Build your own eval set; leaderboards are a guide, not a "
     "guarantee"),
    ("Model sizes",
     "“Open weights” ≠ “runnable locally”: giants are cloud-only",
     "Choose a model class by task and hardware, not by the weight "
     "license"),
]


def build_s38(p):
    """"Let's wrap up": summary table "mechanism -> boundary -> what to
    do" (9 rows, no M-codes, no callback to s01)."""
    s = blank(p)
    slide_title(s, "Let's wrap up: for each mechanism — the boundary "
                   "and what to do about it", size=24, h=0.6)
    tx, ty, tw = 0.55, 1.18, 12.25
    n_rows = len(S38_ROWS) + 1
    gtbl = s.shapes.add_table(n_rows, 3, Inches(tx), Inches(ty),
                              Inches(tw), Inches(4.95))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, w in enumerate([2.55, 5.1, 4.6]):
        tbl.columns[ci].width = Inches(w)
    tbl.rows[0].height = Inches(0.36)
    for ri in range(1, n_rows):
        tbl.rows[ri].height = Inches(0.51)

    def cell(r, c, txt, *, size=10.5, bold=False, color=DEEP, fill=WHITE,
             align=PP_ALIGN.LEFT):
        cl = tbl.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
        cl.margin_left = Inches(0.07); cl.margin_right = Inches(0.05)
        cl.margin_top = Inches(0.01); cl.margin_bottom = Inches(0.01)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cl.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        para.line_spacing = 1.02
        r_ = para.add_run(); r_.text = txt
        r_.font.name = FONT_BODY; r_.font.size = Pt(size)
        r_.font.bold = bold
        r_.font.color.rgb = color

    cell(0, 0, "Mechanism", size=12, bold=True, color=MID)
    cell(0, 1, "Boundary", size=12, bold=True, color=MID)
    cell(0, 2, "What to do", size=12, bold=True, color=MID)
    for ri, (mech, bound, act) in enumerate(S38_ROWS):
        fill = SURFACE if ri % 2 == 0 else WHITE
        cell(ri + 1, 0, mech, size=11, bold=True, color=MID, fill=fill)
        cell(ri + 1, 1, bound, size=10.5, fill=fill)
        cell(ri + 1, 2, act, size=10.5, color=DEEP, fill=fill)
    gold_callout(s, 0.55, 6.32, 12.25, 0.75,
                 "Knowing a tool means knowing its limits. Every "
                 "mechanism works — but not without bound.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s38"))


def build_s39(p):
    """Decision tree "when not an LLM" + escalation ladder."""
    s = blank(p)
    slide_title(s, "When not an LLM — and when not the top LLM",
                size=26, h=0.6)
    text_box(s, 0.55, 1.15, 7.5, 0.4, "When an LLM is not the right "
             "tool:", size=15, bold=True, color=MID)
    branches = [
        ("Classification into fixed categories with thousands of "
         "labeled examples", "classic ML: cheaper, faster, "
         "reproducible — and an LLM is also non-deterministic at T=0"),
        ("Explainability before a regulator", "transparent classic "
         "methods"),
        ("Response < 100 ms (anti-fraud, offline devices)",
         "a specialized small model"),
        ("Exact character-level and arithmetic operations",
         "code, not a model"),
    ]
    yy = 1.62
    for cond, ans in branches:
        ocean_box(s, 0.55, yy, 7.6, 0.92, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, 0.75, yy + 0.05, 7.2, 0.84, [
            {"text": cond, "size": 12, "bold": True, "color": DEEP},
            {"text": "  →  ", "size": 12.5, "bold": True, "color": MID},
            {"text": ans, "size": 12, "color": TEAL, "bold": True},
        ], line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.02
    filled_rect(s, 0.55, yy, 7.6, 0.85, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.4, radius=True, radius_adj=0.1)
    text_runs(s, 0.75, yy + 0.04, 7.2, 0.77, [
        {"text": "Otherwise", "size": 12.5, "bold": True, "color": TEAL},
        {"text": " — language processing, flexible formats, "
                 "multi-step reasoning, generation → ", "size": 12,
         "color": DEEP},
        {"text": "an LLM applies and is often optimal", "size": 12,
         "bold": True, "color": DEEP},
    ], line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
    # Right — escalation ladder
    ocean_box(s, 8.4, 1.15, 4.4, 5.55, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 8.65, 1.3, 3.9, 0.4, "…and not always the top LLM",
             size=15, bold=True, color=DEEP)
    filled_rect(s, 10.3, 1.85, 2.25, 0.8, DEEP, radius=True,
                radius_adj=0.12)
    text_box(s, 10.42, 1.9, 2.0, 0.7, "10% hard cases →\npremium $10/M",
             size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    filled_rect(s, 8.65, 2.8, 3.9, 0.9, MID, radius=True, radius_adj=0.1)
    text_box(s, 8.8, 2.85, 3.6, 0.8, "90% of requests →\na model at "
             "$0.20/M",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.08)
    text_runs(s, 8.65, 3.9, 3.9, 1.2, [
        {"text": "One billion tokens/mo:", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": "$10,000", "size": 14, "bold": True, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
        {"text": " all on premium", "size": 12, "color": DEEP},
        {"text": "vs  ", "size": 12.5, "color": SLATE, "newpara": True,
         "space_before_pt": 6},
        {"text": "$1,180", "size": 16, "bold": True, "color": GOLD},
        {"text": " with routing", "size": 12, "color": DEEP},
    ], line_spacing=1.15)
    # Two Buttons — the hard choice LLM vs. plain code (crop of the top
    # panel with both buttons).
    _place_image_contain(s, ASSETS / "web/twobuttons-llm-vs-code-toponly-en.jpg",
                         8.65, 5.08, 3.9, 1.52)
    speaker_notes(s, load_notes("s39"))


def build_s40(p):
    """Correlation != causation: human vs model + gold callout."""
    s = blank(p)
    slide_title(s, "Attention learns correlation, not causation", size=26,
                h=0.6)
    # Human
    ocean_box(s, 0.55, 1.45, 5.95, 3.55, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.62, 5.35, 0.4, "Humans", size=16.5, bold=True,
             color=MID)
    text_runs(s, 0.85, 2.1, 5.35, 0.85, [
        {"text": "“X happened because Y”", "size": 14, "bold": True,
         "color": DEEP},
        {"text": " — a model of the mechanisms of the world", "size": 13.5,
         "color": DEEP},
    ], line_spacing=1.2)
    hlevels = [("association", "✓"), ("intervention", "✓"),
               ("counterfactual", "✓")]
    yy = 3.0
    for lab, mark in hlevels:
        filled_rect(s, 0.85, yy, 5.35, 0.52, TEAL_TINT, stroke=TEAL,
                    stroke_pt=1.2, radius=True, radius_adj=0.2)
        text_runs(s, 1.05, yy, 5.0, 0.52, [
            {"text": mark + "  ", "size": 13, "bold": True, "color": TEAL},
            {"text": lab, "size": 13, "bold": True, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.62
    # Model
    ocean_box(s, 6.85, 1.45, 5.95, 3.55, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 7.15, 1.62, 5.35, 0.4, "The model (via attention)",
             size=16.5, bold=True, color=DEEP)
    text_runs(s, 7.15, 2.1, 5.35, 0.85, [
        {"text": "“X follows Y in the texts”", "size": 14, "bold": True,
         "color": DEEP},
        {"text": " — “because” for the model is a frequency pattern, "
                 "not a pointer to a mechanism of the world", "size": 12.5,
         "color": DEEP},
    ], line_spacing=1.18)
    mlevels = [
        ("association — strong", "✓", TEAL_TINT, TEAL, TEAL),
        ("intervention — only ones resembling the corpus", "~",
         WHITE, LIGHT, MID),
        ("counterfactual — systematically no", "×", SOFT_GREY, SLATE,
         SLATE),
    ]
    yy = 3.0
    for lab, mark, fill, stroke, mcol in mlevels:
        filled_rect(s, 7.15, yy, 5.35, 0.52, fill, stroke=stroke,
                    stroke_pt=1.2, radius=True, radius_adj=0.2)
        text_runs(s, 7.35, yy, 5.0, 0.52, [
            {"text": mark + "  ", "size": 13, "bold": True, "color": mcol},
            {"text": lab, "size": 11.5, "bold": True, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        yy += 0.62
    # Callout narrowed to ~7.5" — right flank holds xkcd #552 (4.5"
    # wide, 3 panels readable from a projector; 2x variant 918x371 used
    # for sharpness).
    gold_callout(s, 0.55, 5.5, 7.55, 1.35,
                 "Where causal conclusions are expected from the model, "
                 "a human in the loop is an architectural requirement, "
                 "not a polite caveat.", size=15.5, align=PP_ALIGN.CENTER)
    xk_w = 4.5
    xk_h = xk_w * 371 / 918
    _place_image_contain(s, ASSETS / "web/xkcd-552-correlation-2x.png",
                         8.25, 5.38, xk_w, xk_h)
    speaker_notes(s, load_notes("s40"))


def build_s41(p):
    """Bridge to Lecture 3: bridge as a HORIZONTAL STRIP under the
    title, fully readable (arch + cables + towers), not covered by
    cards. 4 concept cards — a compact 2x2 grid below the strip."""
    s = blank(p)
    text_runs(s, 0.55, 0.30, 12.3, 0.55, [
        {"text": "Lecture 3: ", "size": 24, "bold": True, "color": GOLD},
        {"text": "how a model reaches beyond its context", "size": 24,
         "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # Hero strip: bridge cropped to content (1984x784) — reads in full:
    # arch, cables, both towers, end anchors left/right.
    bridge_w = 7.4
    bridge_h = bridge_w * 784 / 1984
    bridge_x = (SLIDE_W_IN - bridge_w) / 2
    bridge_y = 0.90
    add_image(s, ASSETS / "illustrations/s41-bridge-lec3-crop.png",
              x=bridge_x, y=bridge_y, w=bridge_w)
    cards = [
        ("s41-search", "RAG",
         "Semantic search over your own knowledge base → retrieved "
         "fragments go into the context.",
         "Anchor: similarity ≠ relevance.", TEAL),
        ("s41-settings", "Tools / function calling",
         "The model generates a structured call → an external system "
         "executes it.",
         "Anchor: structured outputs guarantee call format.", TEAL),
        ("s41-plug", "MCP",
         "An open protocol for connecting tools.",
         "Anchor: stable prefix → prompt caching, agent economics.",
         GOLD),
        ("s41-refresh-cw", "Agentic loop",
         "Action → observation → correction.",
         "Anchor: the agent reads external content → prompt "
         "injection.", TEAL),
    ]
    card_w, card_h = 5.75, 1.62
    gap_x, gap_y = 0.75, 0.14
    x0 = (SLIDE_W_IN - card_w * 2 - gap_x) / 2
    y0 = bridge_y + bridge_h + 0.14
    for i, (icon, title, body, anchor, acol) in enumerate(cards):
        col, row = i % 2, i // 2
        x = x0 + col * (card_w + gap_x)
        y = y0 + row * (card_h + gap_y)
        ocean_box(s, x, y, card_w, card_h, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.4)
        add_image(s, ASSETS / f"icons/{icon}.png", x=x + 0.18, y=y + 0.14,
                  w=0.34)
        text_box(s, x + 0.62, y + 0.11, card_w - 0.78, 0.36, title,
                 size=13.5, bold=True, color=DEEP)
        text_box(s, x + 0.22, y + 0.50, card_w - 0.44, 0.55, body, size=10.5,
                 color=DEEP, line_spacing=1.08)
        text_box(s, x + 0.22, y + 1.06, card_w - 0.44, 0.50, anchor,
                 size=9.5, italic=True,
                 color=acol if acol == GOLD else TEAL, line_spacing=1.08)
    speaker_notes(s, load_notes("s41"))


def build_s42(p):
    """Q&A minimal (Lec-1 s31 pattern): surface bg, Q&A 140pt, Thank
    you."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, 0.55, 1.7, 12.3, 2.6, "Q&A", size=140, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.55, 4.5, 12.3, 0.8, "Thank you", size=36, italic=True,
             color=MID, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s42"))


# ============================================================
# Build — full deck EN v1.0: 47 slides, 7 sections
# ============================================================
BUILDERS = [
    # Section 0. Introduction
    ("s01", build_s01), ("s02", build_s02), ("s02a", build_s02a),
    ("s03", build_s03), ("s04", build_s04), ("s04b", build_s04b),
    # Section 1. Tokenization
    ("s05a", build_s05a), ("s05", build_s05), ("s06", build_s06),
    ("s08", build_s08), ("s09", build_s09),
    ("s10", build_s10), ("s11", build_s11),
    # Section 2. Embeddings
    ("s12a", build_s12a), ("s12", build_s12), ("s13", build_s13),
    ("s14", build_s14), ("s15", build_s15), ("s17", build_s17),
    # Section 3. Attention (KV-cache right after Q/K/V, role after cache)
    ("s18a", build_s18a), ("s18", build_s18), ("s19", build_s19),
    ("s21", build_s21), ("s22", build_s22), ("s20", build_s20),
    ("s23", build_s23), ("s25", build_s25),
    # Section 4. Sampling and Generation
    ("s26a", build_s26a), ("s26", build_s26), ("s27", build_s27),
    ("s28", build_s28), ("s29", build_s29), ("s30", build_s30),
    ("s31", build_s31), ("s32", build_s32),
    # Section 5. Model Types and Sizes (NEW)
    ("s33a", build_s33a), ("s33", build_s33), ("s34", build_s34),
    ("s36", build_s36), ("s37", build_s37),
    # Section 6. Wrap-up
    ("s35a", build_s35a), ("s35", build_s35), ("s38", build_s38),
    ("s39", build_s39), ("s40", build_s40), ("s41", build_s41),
    ("s42", build_s42),
]


def main():
    p = setup_pres()
    print(f"Building {len(BUILDERS)} slides (EN full deck v1.0)…")
    for sid, fn in BUILDERS:
        try:
            fn(p)
            print(f"  {sid} OK")
        except Exception as e:
            print(f"  {sid} FAIL: {type(e).__name__}: {e}")
            raise
    # EN render is footer-less by convention — no page_number() calls at
    # all (bilingual publish-channel requirement; see module docstring).
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"\nSaved: {OUT}  ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
