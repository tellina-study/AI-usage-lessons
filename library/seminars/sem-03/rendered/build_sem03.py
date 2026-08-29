"""
Build script for Семинар 3 — «Смысл, фокус, случайность: три механизма внутри модели».

Format matched to library/seminars/sem-02 build patterns (same helper
functions, same Ocean Gradient v3 palette, same Q/A split architecture from
the start — no Round-2 retrofit needed here since sem-02's lesson is already
applied: every predict-then-reveal round is TWO slides from slide 1).

Source-of-truth: deck.yaml + slides/*.md.

Canvas: 13.333" x 7.5" (16:9).

Direct python-pptx build (not PowerPoint MCP), per notes/mcp-limitations.md
[#54-1/#54-2/#54-3]: MCP has no list_shapes, format_runs is buggy, no
update_shape_position. Full-rebuild-per-iteration via python-pptx sidesteps
all three. Same choice sem-01/sem-02 made.

Hero images (6-tier acquisition, Tier 2 — Wikimedia directly):
- s01 hero: «T-SNE Embedding of MNIST» (Kyle McDonald, CC BY 2.0) — real
  t-SNE visualization of embedding space, 10 color-coded clusters. Directly
  illustrates the seminar's keystone axis (Block 1): geometric closeness in
  vector space = closeness in meaning.
- s28 hero (closing): «Industrial-robots-21» (Haophuong21, CC BY-SA 4.0) —
  real photo of an industrial robot arm autonomously handling bottles on a
  production line. Bridge to Lecture 3 (agent act -> observe -> correct
  loop; tool calling acting on the world beyond the model's closed cycle).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED Ocean Gradient v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
CODE_BG = RGBColor(0x1B, 0x22, 0x3B)
CODE_FG = RGBColor(0xE3, 0xE9, 0xF2)
GOLD_DARK = RGBColor(0x8A, 0x62, 0x00)  # WCAG-safe dark-gold for text on light bg
                                         # (gold #F0AB00 text fails WCAG AA on
                                         # SURFACE/WHITE — known palette defect,
                                         # see project_ocean_palette_gold_contrast_defect)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons/rendered"
SEM02_ICONS = ROOT.parent / "sem-02/rendered/assets/icons/rendered"
SHOTS = ROOT / "assets/screenshots"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/sem-03.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Courier New"


# ============================================================
# Helpers (adapted from library/seminars/sem-02/rendered/build_sem02.py)
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multipara_box(slide, x, y, w, h, paragraphs, *,
                   anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """Each item in `paragraphs` is a dict of text_box-style kwargs (text/size/...).
    Uses tf.add_paragraph() per line — the ONLY reliable way to force a line break
    (see notes/mcp-limitations.md [#sem01-render-1])."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, cfg in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", align)
        p.line_spacing = cfg.get("line_spacing", 1.15)
        p.space_after = Pt(cfg.get("space_after", 0))
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", FONT_BODY)
        r.font.size = Pt(cfg.get("size", 14))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                 radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def dashed_box(slide, x, y, w, h, *, fill=SURFACE, stroke=GOLD, stroke_pt=1.8,
               radius_pt=12.0, dash="dash"):
    shp = ocean_box(slide, x, y, w, h, fill=fill, stroke=stroke, stroke_pt=stroke_pt,
                     radius_pt=radius_pt)
    ln = shp.line._get_or_add_ln()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    dash_el = etree.SubElement(ln, ns + "prstDash")
    dash_el.set("val", dash)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    path = Path(path)
    if not path.exists():
        print(f"WARNING: missing image {path}")
        return None
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w), height=Inches(h))
    elif w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def icon_path(name, color_hex, size_px):
    """Resolve an icon PNG, preferring sem-03's own asset dir, falling back
    to sem-02's already-rendered 828-icon library (same Lucide set, same
    Ocean recolor — no need to re-download/re-recolor icons that already
    exist for this palette)."""
    local = ICONS / f"{name}-{color_hex}-{size_px}.png"
    if local.exists():
        return local
    shared = SEM02_ICONS / f"{name}-{color_hex}-{size_px}.png"
    return shared


def icon(slide, name, color_hex, size_px, x, y, w_in):
    """Embed a pre-rendered recolored icon PNG (square) at x,y with width w_in (height = width, square icons)."""
    path = icon_path(name, color_hex, size_px)
    return add_image(slide, path, x, y, w=w_in, h=w_in)


def slide_title(slide, text, *, y=0.45, h=1.15, w=12.3, x=0.55, size=28,
                color=DEEP, bold=True, line_spacing=1.15, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=14, bold=True):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.06, w=w - 0.4, h=h - 0.12, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.22)


def speaker_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r"## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)", md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r"\n+---\s*$", "", notes)
    return notes.strip()


def section_tag(slide, x, y, text, *, color=TEAL):
    """Small uppercase section-context label (top-left of body slides).
    NOT a timing marker — section name only, per CLAUDE.md no-timing rule."""
    text_box(slide, x, y, 8.5, 0.32, text=text.upper(), size=11.5, bold=True,
             color=color, align=PP_ALIGN.LEFT)


def vote_hint_bar(slide, x, y, w, h, text, *, fill=SURFACE, stroke=TEAL):
    """Neutral (non-gold) hint bar for Q-slides — voting mechanic only, NO answer."""
    ocean_box(slide, x, y, w, h, fill=fill, stroke=stroke, stroke_pt=1.3)
    icon(slide, "hand", "028090", 64, x + 0.22, y + (h - 0.34) / 2, 0.34)
    text_box(slide, x + 0.75, y, w - 0.95, h, text=text, size=13.5, bold=True,
             color=MID, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# Slide builders — Раздел 1 (s01-s02)
# ============================================================

def build_s01(p):
    """hero_cover — real t-SNE embedding-space visualization (MNIST digit
    clusters, Wikimedia Commons CC BY 2.0), foreshadowing the Block 1
    keystone axis: closeness in vector space = closeness in meaning."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    text_box(s, x=0.55, y=0.62, w=6.35, h=0.45,
             text="ПОСЛЕ ЛЕКЦИИ 2 · СЕМИНАР 3",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.2, w=6.35, h=2.0,
             text="Смысл, фокус, случайность",
             size=38, bold=True, color=DEEP, line_spacing=1.08)
    text_box(s, x=0.55, y=2.95, w=6.35, h=1.15,
             text="Три механизма внутри модели",
             size=20, italic=True, color=MID, line_spacing=1.3)
    text_box(s, x=0.55, y=4.35, w=6.05, h=1.5,
             text="Эмбеддинги, внимание, сэмплинг — три механизма Лекции 2 "
                  "сегодня становятся инструментом предсказания. Вы "
                  "предсказываете, что сделает модель, до того как увидите "
                  "ответ.",
             size=13.5, color=SLATE, line_spacing=1.4)
    text_box(s, x=0.55, y=6.85, w=6.35, h=0.4,
             text="МГТУ им. Н.Э. Баумана", size=12, color=SLATE)

    hero_x, hero_y, hero_w, hero_h = 6.85, 0.0, 6.483, 7.5
    filled_rect(s, hero_x, hero_y, hero_w, hero_h, DEEP)
    img_path = SHOTS / "s01-tsne-embedding-real.png"
    # source is ~1417x1369 (near-square) on white bg — inset on white card
    # for contrast against the dark surround (matches sem-02 s01 pattern).
    # Iter-2 fix: caption now sits ABOVE the white card (in the dark margin)
    # instead of overlapping the gold frame top edge (iter-1 P1 bug — text
    # was unreadable against both gold border and white card).
    pad = 0.5
    white_pad = 0.18
    caption_h = 0.5
    avail_w = hero_w - 2 * pad
    avail_h = hero_h - 2 * pad - 0.55 - caption_h
    card_y = hero_y + pad + 0.25 + caption_h
    filled_rect(s, hero_x + pad - white_pad, card_y - white_pad,
                avail_w + 2 * white_pad, avail_h + 2 * white_pad, WHITE)
    pic = add_image(s, img_path, hero_x + pad, card_y, w=avail_w, h=avail_h)
    if pic is not None:
        gold_frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(hero_x + pad - white_pad), Inches(card_y - white_pad),
            Inches(avail_w + 2 * white_pad), Inches(avail_h + 2 * white_pad))
        gold_frame.fill.background()
        gold_frame.line.color.rgb = GOLD
        gold_frame.line.width = Pt(2.5)
        disable_shadow(gold_frame)
    text_box(s, hero_x + pad, hero_y + pad - 0.08, avail_w, caption_h,
             text="БЛИЗОСТЬ В ПРОСТРАНСТВЕ = БЛИЗОСТЬ СМЫСЛА", size=11.5, bold=True,
             color=GOLD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15)
    text_box(s, hero_x + pad, hero_y + hero_h - pad - 0.1, avail_w, 0.35,
             text="Kyle McDonald · Wikimedia Commons · CC BY 2.0",
             size=9, italic=True, color=RGBColor(0xC8, 0xD2, 0xDF),
             align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Раздел 1 — bridge from Sem-2 (gold callout) + 4-stage pipeline strip,
    stages 2-4 in focus, stage 1 (tokenization) muted."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Раздел 1 · Повтор конвейера")
    slide_title(s, "Четыре этапа конвейера — сегодня фокус на трёх последних", y=0.75, size=25)

    bridge_y = 1.85
    ocean_box(s, 0.55, bridge_y, 12.23, 1.0, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, 0.85, bridge_y + 0.12, 11.6, 0.78,
             text="На Семинаре 2 вы классифицировали AI снаружи. Сегодня — то же "
                  "умение предсказания, но изнутри: эмбеддинги, внимание, сэмплинг "
                  "становятся инструментом.",
             size=14.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    steps_y = bridge_y + 1.3
    steps = [
        ("hash", "Токенизация", "текст → id из словаря", True),
        ("layers", "Эмбеддинг", "id → вектор смысла", False),
        ("scan", "Внимание", "какие токены важны сейчас", False),
        ("shuffle", "Сэмплинг", "распределение → один токен", False),
    ]
    n = len(steps)
    gap = 0.3
    arrow_w = 0.5
    box_w = (12.23 - gap * (n - 1) - arrow_w * (n - 1)) / n
    step_h = 6.55 - steps_y
    cx = 0.55
    for i, (ic, title, sub, muted) in enumerate(steps):
        col = SOFT_GREY if muted else (MID if i == 1 else (TEAL if i == 2 else DEEP))
        stroke = SLATE if muted else LIGHT
        ocean_box(s, cx, steps_y, box_w, step_h, fill=WHITE if not muted else SURFACE, stroke=stroke,
                  stroke_pt=1.2 if muted else 1.7)
        icon_color = "6B7685" if muted else ("065A82" if i == 1 else ("028090" if i == 2 else "21295C"))
        icon(s, ic, icon_color, 96, cx + (box_w - 0.5) / 2, steps_y + 0.3, 0.5)
        text_box(s, cx + 0.08, steps_y + 1.0, box_w - 0.16, 0.4,
                 text=title, size=14.5, bold=True, color=SLATE if muted else DEEP,
                 align=PP_ALIGN.CENTER)
        text_box(s, cx + 0.08, steps_y + 1.45, box_w - 0.16, 0.7,
                 text=sub, size=10.5, italic=True, color=SLATE,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)
        cx += box_w
        if i < n - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                Inches(cx + 0.04), Inches(steps_y + step_h / 2 - 0.14),
                Inches(arrow_w - 0.08), Inches(0.28))
            arr.fill.solid(); arr.fill.fore_color.rgb = TEAL
            arr.line.fill.background()
            disable_shadow(arr)
            cx += arrow_w

    text_box(s, 0.55, 6.75, 12.23, 0.45,
             text="«strawberry» — три токена, поэтому модель плохо считает буквы — этап 1, сегодня не тренируем отдельно",
             size=11.5, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s02"))


# ============================================================
# Slide builders — Блок 1: эмбеддинги (s03-s10)
# ============================================================

def build_s03(p):
    """Раздел 2 intro — recap SSL/HTTPS/React/borsch (5-sentence cosine
    example) + Block 1 section header."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Раздел 2 · Блок 1 — эмбеддинги")
    slide_title(s, "Косинусная близость — угол между векторами смысла", y=0.75, size=26)

    card_y = 1.85
    card_h = 3.5
    ocean_box(s, 0.55, card_y, 12.23, card_h)
    pad = 0.28
    icon(s, "layers", "065A82", 96, 0.55 + pad, card_y + pad, 0.5)
    text_box(s, 0.55 + pad + 0.68, card_y + pad, 6.0, 0.5,
             text="Повтор с Лекции 2", size=16, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("«Как настроить SSL» / «Установка HTTPS-сертификата»", "cosine ≈ 0.85", GOLD_DARK),
        ("«Деплой React-компонента» / «Сборка React-приложения»", "cosine ≈ 0.78", MID),
        ("«Рецепт борща» против любого технического", "cosine ≈ 0.05–0.15", SLATE),
    ]
    ry = card_y + pad + 0.75
    rh = (card_h - pad - 0.75 - pad) / len(rows)
    for text, val, col in rows:
        text_box(s, 0.55 + pad, ry, 8.4, rh, text=text, size=13, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
        text_box(s, 0.55 + pad + 8.5, ry, 3.0, rh, text=val, size=14, bold=True,
                 color=col, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        ry += rh
    text_box(s, 0.55 + pad, card_y + card_h - pad - 0.05, 12.23 - 2 * pad, 0.3,
             text="illustrative — порядок значений устойчив, точные числа зависят от модели",
             size=10.5, italic=True, color=SLATE)

    tag_y = card_y + card_h + 0.35
    text_box(s, 0.55, tag_y, 12.23, 0.5,
             text="Блок 1 — эмбеддинги и семантическая близость", size=19, bold=True, color=DEEP)
    text_box(s, 0.55, tag_y + 0.55, 12.23, 0.5,
             text="Сегодня проверим эту интуицию на новых, более коварных примерах",
             size=13.5, italic=True, color=SLATE)
    speaker_notes(s, load_notes("s03"))


def two_card_pair(p, *, slide_id, section_label, title, cards, hint_text,
                   mode="question", answer_color=None):
    """Shared builder for the omonyms/paraphrases pairs (s04-s07): 2 equal
    Ocean-box cards side by side. `cards` is a list of dicts:
      {"label": "Пара А", "sent1": ..., "sent2": ..., "highlight": bool,
       "cosine": "cosine ≈ 0.1–0.2" (answer mode only),
       "note": "..." (answer mode only)}
    `mode="question"` — no cosine badge, no answer. `mode="answer"` — badge
    + note shown, using `answer_color` (GOLD_DARK-ish tag colour per pair)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, section_label)
    slide_title(s, title, y=0.75, size=25)

    grid_y = 1.95
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.15 if mode == "answer" else 3.55
    for i, c in enumerate(cards):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.26
        text_box(s, cx + pad, grid_y + pad, cw - 2 * pad, 0.4,
                 text=c["label"], size=14, bold=True, color=MID)
        text_box(s, cx + pad, grid_y + pad + 0.5, cw - 2 * pad, 1.5,
                 text=f"«{c['sent1']}»", size=13, italic=True, color=DEEP, line_spacing=1.32)
        text_box(s, cx + pad, grid_y + pad + 1.55, cw - 2 * pad, 1.5,
                 text=f"«{c['sent2']}»", size=13, italic=True, color=DEEP, line_spacing=1.32)
        if mode == "answer":
            badge_y = grid_y + ch - pad - 0.85
            filled_rect(s, cx + pad, badge_y, cw - 2 * pad, 0.5,
                        GOLD_TINT if c.get("gold_badge") else SOFT_GREY,
                        stroke=GOLD if c.get("gold_badge") else SLATE, stroke_pt=1.2,
                        radius=True, radius_adj=0.2)
            text_box(s, cx + pad, badge_y, cw - 2 * pad, 0.5, text=c["cosine"], size=15,
                     bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            text_box(s, cx + pad, badge_y + 0.58, cw - 2 * pad, 0.35,
                     text=c["note"], size=11, italic=True, color=SLATE, line_spacing=1.2)

    hint_y = grid_y + ch + 0.3
    if mode == "question":
        vote_hint_bar(s, 0.55, hint_y, 12.23, 0.55, hint_text)
    else:
        text_box(s, 0.55, hint_y, 12.23, 0.5, text=hint_text, size=13.5, bold=True,
                 color=DEEP, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes(slide_id))


def build_s04(p):
    two_card_pair(p, slide_id="s04", section_label="Блок 1 · Пары А и Б",
        title="«Мышь» и «ключ» — одно слово, два смысла. Близко или далеко?",
        cards=[
            {"label": "Пара А", "sent1": "Кот поймал мышь во дворе",
             "sent2": "Я купил беспроводную мышь для ноутбука"},
            {"label": "Пара Б", "sent1": "Скрипичный ключ в начале нотного стана",
             "sent2": "Забыл ключ от квартиры дома"},
        ],
        hint_text="3 раунда поднятия руки на пару — близко? далеко? не уверен?",
        mode="question")


def build_s05(p):
    two_card_pair(p, slide_id="s05", section_label="Блок 1 · Пары А и Б · ответ",
        title="Слово совпало — смысл разошёлся",
        cards=[
            {"label": "Пара А", "sent1": "Кот поймал мышь во дворе",
             "sent2": "Я купил беспроводную мышь для ноутбука",
             "cosine": "cosine ≈ 0.1–0.2", "note": "мышь-животное vs мышь-устройство",
             "gold_badge": False},
            {"label": "Пара Б", "sent1": "Скрипичный ключ в начале нотного стана",
             "sent2": "Забыл ключ от квартиры дома",
             "cosine": "cosine ≈ 0.1–0.2", "note": "нотный знак vs предмет для замка",
             "gold_badge": False},
        ],
        hint_text="Одинаковые слова — разное значение → низкая близость несмотря на 100% совпадение",
        mode="answer")


def build_s06(p):
    two_card_pair(p, slide_id="s06", section_label="Блок 1 · Пары В и Г",
        title="«Продакшн-сервер» и «бэкенд» — разные слова. Близко или далеко?",
        cards=[
            {"label": "Пара В", "sent1": "Почему упал продакшн-сервер сегодня ночью",
             "sent2": "Причина падения бэкенда прошлой ночью"},
            {"label": "Пара Г", "sent1": "Исправить баг в обработке платежей",
             "sent2": "Починить дефект в модуле оплаты"},
        ],
        hint_text="3 раунда поднятия руки на пару — близко? далеко? не уверен?",
        mode="question")


def build_s07(p):
    two_card_pair(p, slide_id="s07", section_label="Блок 1 · Пары В и Г · ответ",
        title="Слова разные — смысл совпал",
        cards=[
            {"label": "Пара В", "sent1": "Почему упал продакшн-сервер сегодня ночью",
             "sent2": "Причина падения бэкенда прошлой ночью",
             "cosine": "cosine ≈ 0.8+", "note": "«упал»/«падения», «сервер»/«бэкенд»",
             "gold_badge": True},
            {"label": "Пара Г", "sent1": "Исправить баг в обработке платежей",
             "sent2": "Починить дефект в модуле оплаты",
             "cosine": "cosine ≈ 0.8+", "note": "«исправить»/«починить», «баг»/«дефект»",
             "gold_badge": True},
        ],
        hint_text="Эмбеддинг ловит то, что важно для смысла, а не поверхностную форму",
        mode="answer")


_STRAWBERRY_DOCS = [
    "Инструкция по выращиванию клубники на даче",
    "Strawberry fields: a guide to growing your own",
    "Сезонные ягоды средней полосы: сроки созревания",
    "Клубный вечер для сотрудников: программа мероприятия",
    "Рецепт варенья из лесной земляники",
]
_STRAWBERRY_ANSWERS = [
    (True, True, ""),
    (False, True, "перевод через эмбеддинг"),
    (False, True, "родовое понятие"),
    (False, False, "ловушка: общий корень «клуб-»"),
    (False, True, "близкий, не максимальный cosine"),
]


def build_s08(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Блок 1 · Упражнение «клубника»")
    slide_title(s, "Пять документов, запрос «клубника» — кто их найдёт?", y=0.75, size=25)

    intro_y = 1.75
    text_box(s, 0.55, intro_y, 12.23, 0.4,
             text="Запрос «клубника» в корпоративном wiki — тот же пример, что на Лекции 2, теперь на пяти документах",
             size=13, italic=True, color=SLATE)

    table_y = intro_y + 0.5
    table_h = 4.35
    ocean_box(s, 0.55, table_y, 12.23, table_h)
    pad = 0.24
    header_h = 0.4
    text_box(s, 0.55 + pad, table_y + pad, 7.0, header_h, text="Документ", size=12,
             bold=True, color=SLATE)
    text_box(s, 0.55 + pad + 7.3, table_y + pad, 2.15, header_h, text="Полнотекстовый?",
             size=11.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    text_box(s, 0.55 + pad + 9.55, table_y + pad, 2.15, header_h, text="Семантический?",
             size=11.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

    row_y = table_y + pad + header_h + 0.1
    row_h = (table_h - pad - header_h - 0.1 - pad) / 5
    for i, doc in enumerate(_STRAWBERRY_DOCS):
        ry = row_y + i * row_h
        text_box(s, 0.55 + pad, ry, 0.35, row_h, text=str(i + 1), size=13, bold=True,
                 color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 0.55 + pad + 0.42, ry, 6.85, row_h, text=f"«{doc}»", size=12.5,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        icon(s, "help-circle", "6B7685", 64, 0.55 + pad + 7.3 + (2.15 - 0.32) / 2,
             ry + (row_h - 0.32) / 2, 0.32)
        icon(s, "help-circle", "6B7685", 64, 0.55 + pad + 9.55 + (2.15 - 0.32) / 2,
             ry + (row_h - 0.32) / 2, 0.32)
        if i < 4:
            ln = s.shapes.add_connector(1, Inches(0.55 + pad), Inches(ry + row_h),
                                        Inches(0.55 + 12.23 - pad), Inches(ry + row_h))
            ln.line.color.rgb = SOFT_GREY; ln.line.width = Pt(0.75)

    vote_hint_bar(s, 0.55, table_y + table_h + 0.25, 12.23, 0.55,
                  "4 раунда на документ — полнотекстовый? семантический? оба? ни один?")
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Блок 1 · Упражнение «клубника» · ответ")
    slide_title(s, "Общий корень — не то же самое, что общий смысл", y=0.75, size=24)

    table_y = 1.85
    table_h = 4.6
    ocean_box(s, 0.55, table_y, 12.23, table_h)
    pad = 0.22
    header_h = 0.38
    text_box(s, 0.55 + pad, table_y + pad, 6.4, header_h, text="Документ", size=12,
             bold=True, color=SLATE)
    text_box(s, 0.55 + pad + 6.7, table_y + pad, 1.9, header_h, text="Полнотекст",
             size=11, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    text_box(s, 0.55 + pad + 8.65, table_y + pad, 1.9, header_h, text="Семантика",
             size=11, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    text_box(s, 0.55 + pad + 10.6, table_y + pad, 1.3, header_h, text="Заметка",
             size=11, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

    row_y = table_y + pad + header_h + 0.08
    row_h = (table_h - pad - header_h - 0.08 - pad) / 5
    for i, (doc, (ft, sem, note)) in enumerate(zip(_STRAWBERRY_DOCS, _STRAWBERRY_ANSWERS)):
        ry = row_y + i * row_h
        is_trap = (i == 3)
        if is_trap:
            dashed_box(s, 0.55 + 0.05, ry - 0.02, 12.23 - 0.1, row_h + 0.04,
                       fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.6, radius_pt=8.0)
        text_box(s, 0.55 + pad, ry, 0.32, row_h, text=str(i + 1), size=12.5, bold=True,
                 color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 0.55 + pad + 0.38, ry, 6.0, row_h, text=f"«{doc}»", size=11.5,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        ft_icon = "check-circle-2" if ft else "x-circle"
        ft_col = "028090" if ft else "6B7685"
        sem_icon = "check-circle-2" if sem else "x-circle"
        sem_col = "028090" if sem else "6B7685"
        icon(s, ft_icon, ft_col, 64, 0.55 + pad + 6.7 + (1.9 - 0.3) / 2, ry + (row_h - 0.3) / 2, 0.3)
        icon(s, sem_icon, sem_col, 64, 0.55 + pad + 8.65 + (1.9 - 0.3) / 2, ry + (row_h - 0.3) / 2, 0.3)
        if note:
            text_box(s, 0.55 + pad + 10.6, ry, 1.35, row_h, text=note, size=8.5,
                     italic=True, color=GOLD_DARK if is_trap else SLATE,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        if i < 4:
            ln = s.shapes.add_connector(1, Inches(0.55 + pad), Inches(ry + row_h),
                                        Inches(0.55 + 12.23 - pad), Inches(ry + row_h))
            ln.line.color.rgb = SOFT_GREY; ln.line.width = Pt(0.75)

    note_y = table_y + table_h + 0.2
    gold_callout(s, 0.55, note_y, 12.23, 0.55,
                 "Документ 4 — та же логика, что пары А/Б, в обратную сторону: общая часть слова тоже не гарантирует общий смысл",
                 size=12.5)
    speaker_notes(s, load_notes("s09"))


def worked_example_reflection(p, *, slide_id, section_label, title, title_size,
                               example_icon, example_heading, example_body,
                               example_takeaway, question_text, cards):
    """Shared layout for the 3 expanded reflection slides (s10/s20/s25,
    facilitator-guide §2.5/§3.5/§4.4 now 6 min each). Top: wide Ocean box
    worked-example (icon + heading + body + italic takeaway line). Bottom:
    gold-outlined class question + up to 3 compact Ocean-box cards with
    bullet items — same visual grammar as the original two/three-tile
    layout, just compressed to make room for the worked-example above.
    Matches sem-02 `wide_story_bottom_answer` pattern spirit (top block +
    bottom row), adapted for question-only (no colored answer badge, since
    these are open class-discussion prompts, not predict-then-reveal)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, section_label)
    slide_title(s, title, y=0.75, size=title_size)

    ex_y = 1.7
    ex_h = 2.55
    ocean_box(s, 0.55, ex_y, 12.23, ex_h)
    pad = 0.22
    icon(s, example_icon, "065A82", 96, 0.55 + pad, ex_y + pad, 0.42)
    text_box(s, 0.55 + pad + 0.55, ex_y + pad - 0.02, 12.23 - 2 * pad - 0.55, 0.4,
             text=example_heading, size=13.5, bold=True, color=MID)
    text_box(s, 0.55 + pad, ex_y + pad + 0.5, 12.23 - 2 * pad, ex_h - pad - 0.5 - 0.42,
             text=example_body, size=10.8, color=DEEP, line_spacing=1.24)
    text_box(s, 0.55 + pad, ex_y + ex_h - pad - 0.36, 12.23 - 2 * pad, 0.36,
             text=example_takeaway, size=10.5, italic=True, color=TEAL, line_spacing=1.2)

    q_y = ex_y + ex_h + 0.22
    q_h = 0.85
    ocean_box(s, 0.55, q_y, 12.23, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, 0.8, q_y + 0.1, 11.7, q_h - 0.2,
             text=question_text, size=13, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    grid_y = q_y + q_h + 0.2
    grid_h = 7.05 - grid_y
    gap = 0.22
    n = len(cards)
    cw = (12.23 - gap * (n - 1)) / n
    for i, (ic, ctitle, items) in enumerate(cards):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, grid_h)
        cpad = 0.2
        icon(s, ic, "028090", 64, cx + cpad, grid_y + cpad, 0.32)
        text_box(s, cx + cpad + 0.42, grid_y + cpad - 0.03, cw - 2 * cpad - 0.42, 0.44,
                 text=ctitle, size=11.5, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.1)
        iy = grid_y + cpad + 0.48
        item_h = (grid_h - cpad - 0.48 - cpad) / len(items)
        for it in items:
            text_box(s, cx + cpad, iy, 0.2, item_h, text="—", size=11, bold=True, color=LIGHT)
            text_box(s, cx + cpad + 0.22, iy, cw - 2 * cpad - 0.22, item_h,
                     text=it, size=9.8, color=DEEP, line_spacing=1.18)
            iy += item_h
    speaker_notes(s, load_notes(slide_id))


def build_s10(p):
    """Раздел 2.5 — worked-example (юридический поиск по договорам) +
    mandatory not-rhetorical class question, AI-Failure content (6 min,
    facilitator-guide §2.5)."""
    worked_example_reflection(
        p, slide_id="s10", section_label="Блок 1 · Разобранный пример + вопрос всем классом",
        title="Когда полнотекстовый поиск лучше семантического?", title_size=23,
        example_icon="scale",
        example_heading="Разобранный пример: юридический поиск по договорам",
        example_body="Запрос: «пункт об одностороннем расторжении». Полнотекстовый поиск по точной "
                     "фразе находит пункты, где формулировка встречается дословно — критично, потому "
                     "что основание расторжения зафиксировано конкретной формулировкой с юридическими "
                     "последствиями. Семантический поиск дополнительно найдёт релевантный перифраз "
                     "«отказ от договора в одностороннем порядке» — но с не меньшей вероятностью "
                     "подмешает «расторжение по соглашению сторон» (близко по теме, но юридически "
                     "другое основание). Подмена основания расторжения для юриста — не «почти то же "
                     "самое», а содержательная ошибка.",
        example_takeaway="Здесь семантическая близость — источник риска, а не преимущество: задаче "
                          "нужна точность формулировки, а не близость темы.",
        question_text="Каким ещё инженерным задачам, где формулировка/точное значение важнее общей темы, "
                       "нужна именно точность?",
        cards=[
            ("target", "Точность важнее похожести",
             ["Точный поиск по ID / артикулу / номеру заказа",
              "Код ошибки (`NullPointerException`) — семантика может увести к похожим, но не точным совпадениям"]),
            ("arrow-right", "Мост к Лекции 3",
             ["Упражнение «клубника» и семантический поиск — база RAG (Retrieval-Augmented Generation)"]),
        ])


# ============================================================
# Slide builders — Блок 2: внимание (s11-s20) — ЯДРО ЗАНЯТИЯ
# ============================================================

def build_s11(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Раздел 3 · Блок 2 — механизм внимания")
    slide_title(s, "Фонарик в тёмной комнате — внимание распределяет вес между токенами", y=0.75, size=23)

    card_y = 1.85
    card_h = 3.75
    ocean_box(s, 0.55, card_y, 12.23, card_h)
    pad = 0.28
    icon(s, "scan", "065A82", 96, 0.55 + pad, card_y + pad, 0.5)
    text_box(s, 0.55 + pad + 0.68, card_y + pad, 6.5, 0.5,
             text="Повтор с Лекции 2", size=16, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.55 + pad, card_y + pad + 0.65, 12.23 - 2 * pad, 0.8,
             text="Внимание (attention) — фонарик в тёмной комнате: луч направлен на релевантные "
                  "сейчас токены, остальное — на периферии",
             size=13.5, italic=True, color=DEEP, line_spacing=1.3)

    ex_y = card_y + pad + 1.55
    text_box(s, 0.55 + pad, ex_y, 12.23 - 2 * pad, 0.4,
             text="«Кот съел мышь, потому что она была голодна»", size=15, bold=True, color=DEEP)
    weights = [("мышь", 0.5, MID), ("была", 0.28, LIGHT), ("голодна", 0.14, SLATE)]
    wy = ex_y + 0.5
    wx = 0.55 + pad
    for label, weight, col in weights:
        bw = 3.0 * weight * 2
        filled_rect(s, wx, wy, bw, 0.32, col, radius=True, radius_adj=0.5)
        text_box(s, wx + bw + 0.12, wy - 0.02, 2.0, 0.36, text=label, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        wy += 0.42

    tag_y = card_y + card_h + 0.3
    text_box(s, 0.55, tag_y, 12.23, 0.5,
             text="Блок 2 — механизм внимания: угадай, куда «смотрит» модель", size=19, bold=True, color=DEEP)
    text_box(s, 0.55, tag_y + 0.55, 12.23, 0.5,
             text="Сегодня вы предсказываете этот вес сами — до того, как увидите разбор",
             size=13.5, italic=True, color=SLATE)
    speaker_notes(s, load_notes("s11"))


def coref_slide(p, *, slide_id, section_label, title, icon_name, sentence, note,
                 hint_text, mode="question", answer_label=None, weights=None,
                 explanation=None, dashed=False, dual_answer=None):
    """Shared builder for all Block-2 coreference rounds (Winograd x2,
    technical coref, ambiguous coref). `weights` is a list of (label,
    weight_0to1, color) for the answer-mode attention-weight viz.
    `dual_answer` (ambiguous case A-slide only) is a list of 2 dicts
    {"label", "col", "note"} rendered as 2 equal cards instead of 1 badge."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, section_label)
    slide_title(s, title, y=0.75, size=23)

    card_y = 1.95
    # Iter-2 fix (P1 visual-mass): question-mode card enlarged 2.55->3.6 —
    # iter-1 left ~1.7" of dead whitespace below the hint bar (card ended at
    # 4.5, hint bar at 4.8-5.35, nothing until slide bottom 7.05). Taller
    # card + MIDDLE-anchored icon/text fills the vertical space content-
    # fully instead of padding with empty canvas.
    card_h = 3.6 if mode == "question" else (3.0 if not dual_answer else 2.35)
    box_fn = dashed_box if dashed else ocean_box
    box_fn(s, 0.55, card_y, 12.23, card_h)
    pad = 0.3
    icon_y = card_y + pad if mode != "question" else card_y + (card_h - 0.55) / 2 - 0.35
    icon(s, icon_name, "065A82", 96, 0.55 + pad, icon_y, 0.55)
    text_y = card_y + pad if mode != "question" else card_y
    text_h = card_h - 2 * pad if mode != "question" else card_h - (0.4 if note else 0)
    text_box(s, 0.55 + pad + 0.75, text_y, 12.23 - 2 * pad - 0.75, text_h,
             text=f"«{sentence}»", size=18 if mode == "question" else 16.5, italic=True,
             color=DEEP, line_spacing=1.35, anchor=MSO_ANCHOR.MIDDLE)
    if note:
        text_box(s, 0.55 + pad + 0.75, card_y + card_h - pad - 0.32,
                 12.23 - 2 * pad - 0.75, 0.32, text=note, size=11, italic=True, color=SLATE)

    below_y = card_y + card_h + 0.3
    if mode == "question":
        vote_hint_bar(s, 0.55, below_y, 12.23, 0.65, hint_text)
    elif dual_answer:
        gap = 0.3
        cw = (12.23 - gap) / 2
        ch = 7.05 - below_y
        for i, d in enumerate(dual_answer):
            cx = 0.55 + i * (cw + gap)
            ocean_box(s, cx, below_y, cw, ch, stroke=d["col"], stroke_pt=1.8)
            text_box(s, cx, below_y + 0.2, cw, 0.5, text=d["label"], size=19, bold=True,
                     color=d["col"], align=PP_ALIGN.CENTER)
            text_box(s, cx + 0.3, below_y + 0.85, cw - 0.6, ch - 1.5,
                     text=d["note"], size=12.5, color=DEEP, line_spacing=1.35,
                     align=PP_ALIGN.CENTER)
            chip(s, cx + (cw - 1.6) / 2, below_y + ch - 0.55, 1.6, 0.4, "тоже верно",
                 fill=d["col"], size=11)
    else:
        badge_w = 3.4
        filled_rect(s, 0.55, below_y, badge_w, 0.6, GOLD, radius=True, radius_adj=0.25)
        text_box(s, 0.55, below_y, badge_w, 0.6, text=answer_label, size=17, bold=True,
                 color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        weights_block_h = 0.1
        if weights:
            wx = 0.55 + badge_w + 0.4
            wy = below_y + 0.02
            for label, weight, col in weights:
                bw = 5.5 * max(weight, 0.06)
                filled_rect(s, wx, wy, bw, 0.24, col, radius=True, radius_adj=0.5)
                text_box(s, wx + bw + 0.12, wy - 0.03, 2.2, 0.3, text=label, size=11,
                         color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
                wy += 0.32
            weights_block_h = (wy - below_y) + 0.15
        if explanation:
            exp_y = below_y + max(0.75, weights_block_h)
            exp_h = 7.05 - exp_y
            ocean_box(s, 0.55, exp_y, 12.23, exp_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.2)
            text_box(s, 0.75, exp_y + 0.12, 11.8, exp_h - 0.24, text=explanation,
                     size=12.5, color=DEEP, line_spacing=1.32, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes(slide_id))


def build_s12(p):
    coref_slide(p, slide_id="s12", section_label="Блок 2 · пара Winograd Schema, вариант 1",
        title="Трофей не влез в чемодан, потому что он был слишком большой — кто «он»?",
        icon_name="puzzle",
        sentence="Трофей не влез в чемодан, потому что он был слишком большой",
        note="Реальный академический бенчмарк — Levesque, Davis, Morgenstern, The Winograd Schema Challenge (2011/2012)",
        hint_text="3 раунда — он = трофей? он = чемодан? не уверен?", mode="question")


def build_s13(p):
    coref_slide(p, slide_id="s13", section_label="Блок 2 · пара Winograd Schema, вариант 1 · ответ",
        title="«Он» = трофей — большой предмет не влезает в вместилище",
        icon_name="puzzle",
        sentence="Трофей не влез в чемодан, потому что он был слишком большой",
        note=None,
        hint_text="", mode="answer", answer_label="он = трофей",
        weights=[("трофей", 0.75, MID), ("чемодан", 0.15, SLATE)],
        explanation="Трофей большой, поэтому не влез в чемодан. Иллюстрация распределения весов — не настоящий работающий вживую инструмент.")


def build_s14(p):
    coref_slide(p, slide_id="s14", section_label="Блок 2 · пара Winograd Schema, вариант 2",
        title="Трофей не влез в чемодан, потому что он был слишком маленький — кто «он» теперь?",
        icon_name="puzzle",
        sentence="Трофей не влез в чемодан, потому что он был слишком маленький",
        note=None,
        hint_text="3 раунда — он = трофей? он = чемодан? не уверен?", mode="question")


def build_s15(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Блок 2 · пара Winograd Schema, вариант 2 · ответ")
    slide_title(s, "Одна грамматика — разный ответ", y=0.75, size=25)

    card_y = 1.95
    card_h = 2.0
    ocean_box(s, 0.55, card_y, 12.23, card_h)
    pad = 0.3
    icon(s, "puzzle", "065A82", 96, 0.55 + pad, card_y + pad, 0.5)
    text_box(s, 0.55 + pad + 0.7, card_y + pad, 12.23 - 2 * pad - 0.7, card_h - 2 * pad,
             text="«Трофей не влез в чемодан, потому что он был слишком маленький»",
             size=16, italic=True, color=DEEP, line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

    below_y = card_y + card_h + 0.25
    badge_w = 3.4
    filled_rect(s, 0.55, below_y, badge_w, 0.6, GOLD, radius=True, radius_adj=0.25)
    text_box(s, 0.55, below_y, badge_w, 0.6, text="он = чемодан", size=17, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    wx = 0.55 + badge_w + 0.4
    wy = below_y + 0.03
    for label, weight, col in [("чемодан", 0.75, MID), ("трофей", 0.15, SLATE)]:
        bw = 5.5 * weight
        filled_rect(s, wx, wy, bw, 0.22, col, radius=True, radius_adj=0.5)
        text_box(s, wx + bw + 0.12, wy - 0.05, 2.0, 0.3,
                 text=label, size=11, color=DEEP)
        wy += 0.3

    call_y = below_y + 0.85
    call_h = 7.05 - call_y
    gold_callout(s, 0.55, call_y, 12.23, call_h,
                 "Одна и та же грамматика в обоих предложениях — а ответ разный. Для правильного "
                 "ответа нужны мировые знания, не только синтаксический разбор.", size=14.5)
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    coref_slide(p, slide_id="s16", section_label="Блок 2 · Техническая кореференция",
        title="Менеджер поругал стажёра, потому что он не проверил тесты — кто «он»?",
        icon_name="git-branch",
        sentence="Менеджер поругал стажёра, потому что он не проверил тесты перед коммитом",
        note=None,
        hint_text="3 раунда — он = менеджер? он = стажёр? не уверен?", mode="question")


def build_s17(p):
    coref_slide(p, slide_id="s17", section_label="Блок 2 · Техническая кореференция · ответ",
        title="«Он» = стажёр — вес внимания на «стажёр» и «коммит»",
        icon_name="git-branch",
        sentence="Менеджер поругал стажёра, потому что он не проверил тесты перед коммитом",
        note=None,
        hint_text="", mode="answer", answer_label="он = стажёр",
        weights=[("стажёр", 0.55, MID), ("коммит", 0.3, TEAL), ("менеджер", 0.12, SLATE)],
        explanation="Сущность, которую ругают, связана с конкретным упущенным действием — типичная обязанность того, кто делает коммит.")


def build_s18(p):
    coref_slide(p, slide_id="s18", section_label="Блок 2 · Пограничный случай",
        title="Пилот сказал диспетчеру, что он ошибся в расчётах — кто «он»?",
        icon_name="help-circle",
        sentence="Пилот сказал диспетчеру, что он ошибся в расчётах",
        note="Оба существительных мужского рода — согласование рода не подсказывает ответ",
        hint_text="3 раунда — он = пилот? он = диспетчер? не уверен?", mode="question",
        dashed=True)


def build_s19(p):
    coref_slide(p, slide_id="s19", section_label="Блок 2 · Пограничный случай · ответ",
        title="Оба прочтения разумны",
        icon_name="help-circle",
        sentence="Пилот сказал диспетчеру, что он ошибся в расчётах",
        note=None,
        hint_text="", mode="answer", dashed=True,
        dual_answer=[
            {"label": "он = пилот", "col": LIGHT, "note": "пилот признал свою ошибку диспетчеру"},
            {"label": "он = диспетчер", "col": TEAL, "note": "пилот указал диспетчеру на его ошибку"},
        ])


def build_s20(p):
    """Раздел 3.5 — worked-example (документированный провал WinoBias) +
    closing reflection, Pearl causality-levels connection (6 min,
    facilitator-guide §3.5)."""
    worked_example_reflection(
        p, slide_id="s20", section_label="Блок 2 · Разобранный пример + рефлексия",
        title="Уверенный ответ — это понимание или угадывание?", title_size=23,
        example_icon="alert-triangle",
        example_heading="Документированный провал: WinoBias",
        example_body="Zhao, Wang, Yatskar, Ordonez, Chang (2018), NAACL — отдельный от Winograd "
                     "Schema бенчмарк именно на систематическую ошибку. WinoBias показал: модели "
                     "разрешения кореференции систематически резолвят неоднозначное местоимение через "
                     "профессионально-гендерные стереотипы обучающего корпуса — например, чаще "
                     "связывают «врач» с мужским местоимением, а «медсестра» — с женским, даже когда "
                     "контекст и грамматика этого не требуют.",
        example_takeaway="Прямое следствие примера «пилот/диспетчер»: внимание статистически смотрит "
                          "на корпус, и если корпус несёт социальный перекос (bias — тот же термин, "
                          "что на Семинаре 2), модель воспроизводит его уверенно, а не осторожно.",
        question_text="Если модель уверенно называет референт «он» — значит ли это, что она «понимает» "
                       "предложение, как человек? Чем понимание отличается от угадывания?",
        cards=[
            ("network", "Корреляция", ["Внимание — статистика совместной встречаемости токенов"]),
            ("compass", "Причинность", ["Человек опирается на модель мира — уровни Перла, Лекция 1"]),
            ("help-circle", "Пограничный случай", ["Модель угадывает не хуже человека там, где сам человек не уверен"]),
        ])


# ============================================================
# Slide builders — Блок 3: сэмплинг (s21-s25)
# ============================================================

def build_s21(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Раздел 4 · Блок 3 — сэмплинг")
    slide_title(s, "Четыре ручки API — предсказуемость против разнообразия", y=0.75, size=25)

    card_y = 1.85
    card_h = 3.6
    ocean_box(s, 0.55, card_y, 12.23, card_h)
    pad = 0.28
    icon(s, "gauge", "065A82", 96, 0.55 + pad, card_y + pad, 0.5)
    text_box(s, 0.55 + pad + 0.68, card_y + pad, 6.0, 0.5,
             text="Повтор с Лекции 2", size=16, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("target", "Классификация", "T = 0"),
        ("code", "Кодогенерация", "T ≈ 0.2–0.3"),
        ("message-square", "Чат-объяснение", "T ≈ 0.7"),
        ("sparkles", "Творческое письмо", "T ≈ 0.9–1.2"),
    ]
    ry = card_y + pad + 0.7
    rh = (card_h - pad - 0.7 - pad) / len(rows)
    for ic, label, val in rows:
        icon(s, ic, "028090", 64, 0.55 + pad, ry + (rh - 0.3) / 2, 0.3)
        text_box(s, 0.55 + pad + 0.44, ry, 4.5, rh, text=label, size=13.5, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        val_x = 0.55 + pad + 5.2
        bar_max = 5.8
        frac = {"T = 0": 0.05, "T ≈ 0.2–0.3": 0.25, "T ≈ 0.7": 0.6, "T ≈ 0.9–1.2": 0.9}[val]
        filled_rect(s, val_x, ry + rh / 2 - 0.08, bar_max * frac, 0.16, GOLD if frac > 0.8 else MID,
                    radius=True, radius_adj=0.5)
        text_box(s, val_x + bar_max + 0.15, ry, 2.0, rh, text=val, size=12, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE)
        ry += rh

    tag_y = card_y + card_h + 0.3
    text_box(s, 0.55, tag_y, 12.23, 0.5,
             text="Блок 3 — сэмплинг: температура и параметры", size=19, bold=True, color=DEEP)
    text_box(s, 0.55, tag_y + 0.55, 12.23, 0.5,
             text="Сегодня вы проверите эту логику на собственном эксперименте",
             size=13.5, italic=True, color=SLATE)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """Раздел 4.2 — homework debrief prompt, no facilitator instructions on slide."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Блок 3 · Разбор домашнего задания")
    slide_title(s, "T=0 / T=0.7 / T=1.5 — что изменилось в вашем эксперименте?", y=0.75, size=23)

    tag_y = 1.75
    ocean_box(s, 0.55, tag_y, 12.23, 0.55, fill=SURFACE, stroke=TEAL, stroke_pt=1.2)
    text_box(s, 0.85, tag_y, 11.6, 0.55,
             text="Домашнее задание из Лекции 2",
             size=13, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)

    grid_y = tag_y + 0.75
    gap = 0.26
    cw = (12.23 - 2 * gap) / 3
    ch = 1.85
    cols = [
        ("target", "T = 0", "детерминированный argmax", "065A82"),
        ("shuffle", "T = 0.7", "естественная вариативность", "028090"),
        ("zap", "T = 1.5", "на грани распада связности", "8A6200"),
    ]
    for i, (ic, title, sub, col_hex) in enumerate(cols):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        icon(s, ic, col_hex, 96, cx + (cw - 0.44) / 2, grid_y + 0.2, 0.44)
        text_box(s, cx, grid_y + 0.75, cw, 0.4, text=title, size=17, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER)
        text_box(s, cx + 0.15, grid_y + 1.2, cw - 0.3, 0.55, text=sub, size=11, italic=True,
                 color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.15)

    q_y = grid_y + ch + 0.3
    q_h = 7.05 - q_y
    ocean_box(s, 0.55, q_y, 12.23, q_h)
    pad = 0.24
    questions = [
        "Что изменилось от T=0 к T=1.5?",
        "Что осталось одинаковым во всех трёх запусках?",
        "Какую T вы выбрали для боевого режима и почему?",
    ]
    qy = q_y + pad
    qh = (q_h - 2 * pad) / len(questions)
    for i, q in enumerate(questions):
        text_box(s, 0.55 + pad, qy, 0.4, qh, text=f"{i+1}.", size=14, bold=True, color=LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 0.55 + pad + 0.45, qy, 12.23 - 2 * pad - 0.45, qh, text=q, size=14.5,
                 bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        qy += qh
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Блок 3 · Угадай температуру")
    slide_title(s, "Три ответа на один промпт — какой температуре какой соответствует?", y=0.75, size=21)

    tag_y = 1.65
    ocean_box(s, 0.55, tag_y, 12.23, 0.5, fill=SOFT_GREY, stroke=SLATE, stroke_pt=1.0)
    text_box(s, 0.85, tag_y, 11.6, 0.5,
             text="Три ответа ниже написаны автором методического комплекта для иллюстрации — это не реальный вызов модели",
             size=11, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    prompt_y = tag_y + 0.6
    text_box(s, 0.55, prompt_y, 12.23, 0.4,
             text='Промпт: «Напиши короткое приветствие для email-рассылки студентам о начале нового курса»',
             size=12, italic=True, color=SLATE)

    grid_y = prompt_y + 0.5
    gap = 0.24
    cw = (12.23 - 2 * gap) / 3
    ch = 3.4
    samples = [
        ("А", "Уважаемые студенты! Рады сообщить вам о начале нового курса. Занятия "
              "начинаются согласно расписанию. Желаем успехов в обучении!"),
        ("Б", "Привет! Новый курс стартует уже на следующей неделе — будет практика, "
              "разбор реальных кейсов и, чего уж скрывать, немного домашних заданий. "
              "Ждём вас на первом занятии."),
        ("В", "Курс! Начинается! Проверьте расписание — а может, и не проверяйте, "
              "жизнь непредсказуема, как и первое занятие, где кофе важнее конспекта, "
              "но конспект тоже важен, наверное."),
    ]
    for i, (letter, text) in enumerate(samples):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.2
        text_box(s, cx + pad, grid_y + pad - 0.02, cw - 2 * pad, 0.35,
                 text=f"Ответ {letter}", size=13.5, bold=True, color=MID)
        text_box(s, cx + pad, grid_y + pad + 0.4, cw - 2 * pad, ch - 0.4 - 2 * pad,
                 text=text, size=11, color=DEEP, line_spacing=1.3)

    vote_hint_bar(s, 0.55, grid_y + ch + 0.25, 12.23, 0.55,
                  "3 раунда — для каждого ответа: низкая / средняя / высокая температура?")
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Блок 3 · Угадай температуру · ответ")
    slide_title(s, "Повторяемость → вариативность → распад связности", y=0.75, size=22)

    grid_y = 1.75
    gap = 0.24
    cw = (12.23 - 2 * gap) / 3
    ch = 3.9
    samples = [
        ("А", "Уважаемые студенты! Рады сообщить вам о начале нового курса. Занятия "
              "начинаются согласно расписанию. Желаем успехов в обучении!",
         "T ≈ 0", "полная предсказуемость, текст почти не меняется при повторной генерации", LIGHT),
        ("Б", "Привет! Новый курс стартует уже на следующей неделе — будет практика, "
              "разбор реальных кейсов и, чего уж скрывать, немного домашних заданий. "
              "Ждём вас на первом занятии.",
         "T ≈ 0.7", "конкретные детали, живая интонация, связный текст без странностей", TEAL),
        ("В", "Курс! Начинается! Проверьте расписание — а может, и не проверяйте, "
              "жизнь непредсказуема, как и первое занятие, где кофе важнее конспекта, "
              "но конспект тоже важен, наверное.",
         "T ≈ 1.5–2.0", "неожиданные обороты, резкие смены тона, на грани распада связности", GOLD_DARK),
    ]
    for i, (letter, text, temp, feature, col) in enumerate(samples):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.2
        text_box(s, cx + pad, grid_y + pad - 0.02, cw - 2 * pad, 0.35,
                 text=f"Ответ {letter}", size=13, bold=True, color=MID)
        text_box(s, cx + pad, grid_y + pad + 0.38, cw - 2 * pad, 1.75,
                 text=text, size=10, color=DEEP, line_spacing=1.25)
        badge_y = grid_y + pad + 2.2
        filled_rect(s, cx + pad, badge_y, cw - 2 * pad, 0.4, col, radius=True, radius_adj=0.3)
        text_box(s, cx + pad, badge_y, cw - 2 * pad, 0.4, text=temp, size=13, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + pad, badge_y + 0.48, cw - 2 * pad, ch - pad - 0.48 - (badge_y - grid_y),
                 text=feature, size=10, italic=True, color=SLATE, line_spacing=1.25)

    text_box(s, 0.55, grid_y + ch + 0.2, 12.23, 0.4,
             text="Признак — не подпись, а наблюдаемое поведение текста", size=13, bold=True,
             color=GOLD_DARK, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s24"))


def build_s25(p):
    """Раздел 4.4 — worked-example (банковский чат-бот, regex+Luhn вместо
    LLM-сэмплинга) + closing reflection (6 min, facilitator-guide §4.4)."""
    worked_example_reflection(
        p, slide_id="s25", section_label="Блок 3 · Разобранный пример + рефлексия",
        title="Когда стохастичность — риск, а не удобство?", title_size=23,
        example_icon="shield-check",
        example_heading="Разобранный пример: когда НЕ LLM вообще",
        example_body="Банковский чат-бот использует LLM для диалога, но для детекции номеров "
                     "банковских карт в переписке (чтобы замаскировать их перед логированием) — не "
                     "полагается на LLM ни при каком значении температуры. Используется regex по "
                     "формату номера карты + Luhn-алгоритм проверки контрольной суммы. Причина: нужна "
                     "гарантия срабатывания на 100% известных случаев, а не вероятностное соответствие "
                     "— даже при T=0 LLM может пропустить нестандартно отформатированный номер или "
                     "ложно среагировать на похожую последовательность цифр, а regex — нет.",
        example_takeaway="Тот же принцип, что в Блоке 1 с полнотекстовым поиском: если задаче нужна "
                          "гарантия, а не вероятностное приближение — классический инструмент с "
                          "фиксированной грамматикой правильнее LLM при любом параметре сэмплинга.",
        question_text="В каких сценариях вашей будущей работы стохастичность ответа (T>0) — это риск, "
                       "а не удобство?",
        cards=[
            ("scale", "Где нужна воспроизводимость",
             ["Классификация с аудитом — почему в прошлый раз результат был другим",
              "Тестирование ПО — нужен детерминированный оракул для сравнения"]),
            ("alert-triangle", "Оговорка про T=0",
             ["Даже T=0 не даёт стопроцентной побитовой детерминированности",
              "T=0 — инженерное приближение, не абсолютная гарантия"]),
        ])


# ============================================================
# Slide builders — Раздел 5-6 (s26-s28)
# ============================================================

def build_s26(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Раздел 5 · Мостик к Лекции 3")
    slide_title(s, "Сегодня — три механизма поодиночке. Дальше — модель выходит за пределы цикла", y=0.75, size=21)

    card_y = 1.95
    card_h = 2.4
    ocean_box(s, 0.55, card_y, 12.23, card_h)
    pad = 0.3
    text_box(s, 0.55 + pad, card_y + pad, 12.23 - 2 * pad, card_h - 2 * pad,
             text="Сегодня мы натренировали три механизма внутри модели поодиночке: "
                  "эмбеддинги — для поиска смысла, внимание — для фокуса, сэмплинг — для "
                  "выбора токена. На Лекции 3 модель выходит за пределы этого замкнутого "
                  "цикла — ищет во внешней базе знаний, вызывает инструменты, действует в цикле агента.",
             size=14, italic=True, color=DEEP, line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)

    tiles_y = card_y + card_h + 0.35
    tiles_h = 7.05 - tiles_y
    gap = 0.24
    n = 4
    tw = (12.23 - gap * (n - 1)) / n
    tags = [
        ("search", "RAG"),
        ("puzzle", "Function calling"),
        ("network", "MCP"),
        ("rotate-ccw", "Agent loop"),
    ]
    for i, (ic, label) in enumerate(tags):
        cx = 0.55 + i * (tw + gap)
        filled_rect(s, cx, tiles_y, tw, tiles_h, SURFACE, stroke=TEAL, stroke_pt=1.3, radius=True, radius_adj=0.1)
        icon(s, ic, "028090", 96, cx + (tw - 0.5) / 2, tiles_y + 0.25, 0.5)
        text_box(s, cx + 0.1, tiles_y + tiles_h - 0.55, tw - 0.2, 0.4, text=label, size=13,
                 bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    section_tag(s, 0.55, 0.4, "Раздел 6 · Памятка на вынос")
    slide_title(s, "Три механизма внутри модели — что помнить", y=0.75, size=25)

    grid_y = 1.85
    grid_h = 4.75
    gap = 0.26
    cw = (12.23 - 2 * gap) / 3
    cols = [
        ("layers", "Эмбеддинги = близость по смыслу, не по буквам", MID, [
            "Одинаковые слова могут означать разное — низкая близость несмотря на лексическое совпадение",
            "Разные слова могут означать одно и то же — высокая близость несмотря на отсутствие общих слов",
            "Полнотекстовый поиск иногда точнее семантического: точные ID, юридические цитаты, коды ошибок",
        ]),
        ("scan", "Внимание = статистическая корреляция, не понимание причинности", TEAL, [
            "Модель предсказывает референт через статистику совместной встречаемости токенов",
            "Иногда однозначного ответа не существует — честный результат тогда распределение вероятностей",
        ]),
        ("shuffle", "Сэмплинг = баланс предсказуемость ↔ разнообразие", DEEP, [
            "T=0 — для задач, где важна воспроизводимость, но без стопроцентной гарантии",
            "T>0.7 — для задач, где ценно разнообразие",
            "Правильная T — результат эмпирической калибровки, не константа",
        ]),
    ]
    for i, (ic, title, col, items) in enumerate(cols):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, grid_h)
        pad = 0.22
        filled_rect(s, cx + pad, grid_y + pad, 0.55, 0.55, col, radius=True, radius_adj=0.3)
        icon(s, ic, "FFFFFF", 64, cx + pad + 0.11, grid_y + pad + 0.11, 0.33)
        text_box(s, cx + pad, grid_y + pad + 0.68, cw - 2 * pad, 0.9,
                 text=title, size=12.5, bold=True, color=DEEP, line_spacing=1.22)
        iy = grid_y + pad + 1.68
        item_h = (grid_h - pad - 1.68 - pad) / len(items)
        for it in items:
            text_box(s, cx + pad, iy, 0.2, item_h, text="—", size=11, bold=True, color=LIGHT)
            text_box(s, cx + pad + 0.24, iy, cw - 2 * pad - 0.24, item_h,
                     text=it, size=9.8, color=SLATE, line_spacing=1.22)
            iy += item_h

    text_box(s, 0.55, grid_y + grid_h + 0.2, 12.23, 0.4,
             text="Всё, что здесь написано, вы уже видели сегодня — рабочая версия на вынос",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """hero_closing — real industrial robot arm photo (Wikimedia Commons
    CC BY-SA 4.0), bridging to Lecture 3's agent act->observe->correct loop
    theme."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    hero_w = 6.283
    img_path = SHOTS / "s-closing-industrial-robots-real.jpg"
    filled_rect(s, 0, 0, hero_w, 7.5, DEEP)
    pad = 0.4
    avail_w = hero_w - 2 * pad
    avail_h = 7.5 - 2 * pad - 0.5
    pic = add_image(s, img_path, pad, pad + 0.25, w=avail_w, h=avail_h)
    if pic is not None:
        gframe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pic.left, pic.top, pic.width, pic.height)
        gframe.fill.background()
        gframe.line.color.rgb = GOLD
        gframe.line.width = Pt(2.5)
        disable_shadow(gframe)
    text_box(s, pad, pad - 0.05, avail_w, 0.3, text="ДЕЙСТВИЕ ВО ВНЕШНЕМ МИРЕ", size=11, bold=True,
             color=GOLD, align=PP_ALIGN.LEFT)
    text_box(s, pad, 7.5 - pad - 0.1, avail_w, 0.3,
             text="Haophuong21 · Wikimedia Commons · CC BY-SA 4.0",
             size=8.5, italic=True, color=RGBColor(0xC8, 0xD2, 0xDF))

    rx = hero_w + 0.55
    rw = 13.333 - rx - 0.55
    text_box(s, rx, 1.35, rw, 0.5, text="СПАСИБО", size=14, bold=True, color=TEAL)
    text_box(s, rx, 1.85, rw, 2.1,
             text="От механизмов внутри — к тому, что снаружи",
             size=26, bold=True, color=DEEP, line_spacing=1.15)
    text_box(s, rx, 3.9, rw, 1.8,
             text="Сегодня вы натренировали предсказание эмбеддингов, внимания и "
                  "сэмплинга. На Лекции 3 — как модель выходит за пределы этого цикла "
                  "и действует в мире через RAG, инструменты, MCP, цикл агента.",
             size=14, color=MID, line_spacing=1.4)
    chip(s, rx, 5.95, 3.3, 0.55, "Лекция 3  →  далее", fill=DEEP, size=14)
    speaker_notes(s, load_notes("s28"))


# ============================================================
# Main
# ============================================================
BUILDERS = [
    build_s01, build_s02, build_s03, build_s04, build_s05,
    build_s06, build_s07, build_s08, build_s09, build_s10,
    build_s11, build_s12, build_s13, build_s14, build_s15,
    build_s16, build_s17, build_s18, build_s19, build_s20,
    build_s21, build_s22, build_s23, build_s24, build_s25,
    build_s26, build_s27, build_s28,
]


def main():
    p = setup_pres()
    for fn in BUILDERS:
        fn(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved {OUT} ({len(BUILDERS)} slides)")


if __name__ == "__main__":
    main()
