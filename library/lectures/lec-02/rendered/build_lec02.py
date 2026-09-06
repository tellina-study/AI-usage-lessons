"""
Лекция 2 v3.0 (issue #183, rework round 2) — полный дек, 47 слайдов.

Source-of-truth: deck.yaml v3.0 + slides/*.md (v3 sources).
Порядок v3.0: s01 s02 s02a s03 s04 s04b | s05a s05 s06 s08 s09 s10 s11 |
  s12a s12 s13 s14 s15 s17 | s18a s18 s19 s21 s22 s20 s23 s25 |
  s26a s26 s27 s28 s29 s30 s31 s32 | s33a s33 s34 s36 s37 |
  s35a s35 s38 s39 s40 s41 s42.

v3.0 key changes vs v2.0 (spec: notes/lecture-2-review/final/rework-round2.md):
- s07 (chat-шаблоны) УДАЛЁН — тема перенесена в Лекцию 3.
- Перестановка Раздела 3: s18 → s19 → s21 (KV-cache) → s22 → s20 → s23 → s25.
- НОВЫЙ Раздел 5 «Виды и размеры моделей»: s33a (divider) + s33 (NEW
  классификация по размеру) + s34/s36/s37 (перенос). Финал = Раздел 6.
- Арка чек-листа снята: s01 — двухпанельный мем T=0; s02a — 7 карточек
  без M-чипов; s04 — 7 промисов; s38 — «Подведём итоги» (таблица
  «механизм → граница → что делать»).
- Пример Раздела 3: «Кот съел мышь, потому что ОН был голоден» —
  внимание от «он» к «Кот» (согласование по роду) на s18/s19/s20.
- Мем-иллюстрации на дивайдерах (assets/illustrations, gen_assets_v3.py).

Pipeline: python-pptx direct (канонический production-путь по
notes/mcp-limitations.md [#71-1] — full rebuild each iteration).
Palette LOCKED: Ocean #21295C/#065A82/#1C7293 + teal #028090 + gold #F0AB00.
Canvas 13.333"x7.5" (16:9).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED, mirrored from deck.yaml v2.0.1) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
SLATE_PN  = RGBColor(0x5B, 0x66, 0x78)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered/assets"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-02.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ============================================================
# Helpers (ported from v1.8 builder; refs system dropped)
# ============================================================
def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall(A_NS + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, A_NS + "effectLst")


def set_fill_alpha(shp, opacity_pct):
    """Полупрозрачная заливка: добавляет <a:alpha> в solidFill/srgbClr
    (opacity_pct: 0..100, 100 = непрозрачный). v2.1 — постер s01."""
    spPr = shp._element.spPr
    sf = spPr.find(A_NS + "solidFill")
    if sf is None:
        return
    srgb = sf.find(A_NS + "srgbClr")
    if srgb is None:
        return
    for el in srgb.findall(A_NS + "alpha"):
        srgb.remove(el)
    etree.SubElement(srgb, A_NS + "alpha").set(
        "val", str(int(opacity_pct * 1000)))


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = line_spacing
    r = para.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    """runs: list of dicts {text, size, bold, italic, color, font, newpara,
    align, line_spacing, space_after_pt}."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            para = tf.add_paragraph()
            para.alignment = cfg.get("align", align)
            para.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before_pt"):
                para.space_before = Pt(cfg["space_before_pt"])
        r = para.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=14, bold=True, stroke_pt=1.2, font=FONT_BODY):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    """Prefer only one of w/h — both stretches non-proportionally
    (notes/mcp-limitations.md #73-render-1)."""
    if not Path(path).exists():
        return
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w))
    elif h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.5, h=1.0, w=12.3, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.18,
                align=PP_ALIGN.LEFT):
    text_box(slide, x, y, w, h, text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                radius=True, radius_adj=0.12)
    text_box(slide, x + 0.2, y + 0.06, w - 0.4, h - 0.12, text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.22)


def right_arrow(slide, x, y, w=0.6, h=0.4, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def down_arrow(slide, x, y, w=0.4, h=0.45, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def line_arrow(slide, x1, y1, x2, y2, *, color=MID, w_pt=2.0, dash=None):
    """Прямая линия-стрелка (tailEnd triangle) — для веерных стрелок s20 и
    пунктирных коннекторов s35."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w_pt)
    if dash is not None:
        conn.line.dash_style = dash
    ln = conn.line._get_or_add_ln()
    end = ln.makeelement(A_NS + "tailEnd",
                         {"type": "triangle", "w": "med", "len": "med"})
    ln.append(end)
    disable_shadow(conn)
    return conn


def plain_line(slide, x1, y1, x2, y2, *, color=LIGHT, w_pt=1.5, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w_pt)
    if dash is not None:
        conn.line.dash_style = dash
    disable_shadow(conn)
    return conn


def left_arrow(slide, x, y, w=0.6, h=0.4, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def outline_big_text(slide, x, y, w, h, text, *, size=200, stroke=GOLD,
                     stroke_pt=2.6, align=PP_ALIGN.LEFT):
    """Decorative outline numeral: white fill + gold stroke via a:ln XML."""
    tb = text_box(slide, x, y, w, h, text, size=size, bold=True,
                  color=WHITE, align=align, line_spacing=0.95)
    r = tb.text_frame.paragraphs[0].runs[0]
    rPr = r._r.get_or_add_rPr()
    ln = rPr.makeelement(A_NS + "ln", {"w": str(int(stroke_pt * 12700))})
    sf = etree.SubElement(ln, A_NS + "solidFill")
    clr = etree.SubElement(sf, A_NS + "srgbClr")
    clr.set("val", "F0AB00" if stroke == GOLD else str(stroke))
    rPr.insert(0, ln)
    return tb


def speaker_notes(slide, text):
    """Readable paragraphs: blank-line-separated blocks each become one
    notes paragraph; single newlines collapse to spaces."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    blocks = [b.strip() for b in re.split(r'\n\s*\n', text.strip()) if b.strip()]
    if not blocks:
        blocks = [""]
    for i, block in enumerate(blocks):
        one = re.sub(r'\s*\n\s*', ' ', block)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = one


def load_notes(slide_id):
    files = sorted(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\Z)', md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes).strip()
    # v3.0: strip-safe внутренние маркеры источников ([VFY-day-of] /
    # [FACT-CHECK: …]) не выпекаются в speaker notes pptx (pre-gate grep = 0)
    notes = re.sub(r'\s*\[(?:VFY-day-of|FACT-CHECK[^\]]*)\]', '', notes)
    return notes


def page_number(slide, n, *, color=SLATE_PN):
    tb = slide.shapes.add_textbox(Inches(12.15), Inches(7.16), Inches(1.1),
                                  Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    r = para.add_run(); r.text = str(n)
    r.font.name = FONT_BODY; r.font.size = Pt(10)
    r.font.italic = True; r.font.color.rgb = color
    return tb


# ============================================================
# v2.0 pipeline progress bar (dividers) — 7 стадий конвейера s04b.
# Активная стадия(и) gold. Без «вы здесь», без минут.
# ============================================================
PIPE_STAGES = ["Текст", "Токены", "Векторы", "LLM",
               "Распределение", "Токен", "Текст"]


def pipeline_bar(slide, active, *, y=6.62, bar_h=0.48, passed=frozenset(),
                 muted_gold=False):
    """active: int|set — gold-подсветка. v3.0: passed — уже пройденные
    стадии (teal-tint, briefs s18a/s26a); muted_gold=True — ВСЯ лента в
    приглушённом gold (s33a: раздел логически «над» конвейером, не
    отдельная стадия; визуально слабее полного gold s35a)."""
    if isinstance(active, int):
        active = {active}
    n = len(PIPE_STAGES)
    arrow_w = 0.24
    total_w = 12.3
    cell_w = (total_w - arrow_w * (n - 1)) / n
    start_x = (SLIDE_W_IN - total_w) / 2.0
    for i, label in enumerate(PIPE_STAGES):
        x = start_x + i * (cell_w + arrow_w)
        is_act = (not muted_gold) and (i in active)
        if muted_gold:
            fill, color = GOLD_TINT, DEEP
        elif is_act:
            fill, color = GOLD, DEEP
        elif i in passed:
            fill, color = TEAL_TINT, TEAL
        else:
            fill, color = SOFT_GREY, SLATE
        filled_rect(slide, x, y, cell_w, bar_h, fill, radius=True,
                    radius_adj=0.28)
        sz = 11 if len(label) < 12 else 9.5
        text_box(slide, x - 0.05, y + 0.06, cell_w + 0.1, bar_h - 0.12,
                 label, size=sz, bold=is_act, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            ax = x + cell_w + 0.03
            right_arrow(slide, ax, y + bar_h / 2 - 0.07, w=arrow_w - 0.06,
                        h=0.14, fill=LIGHT)


def section_divider(p, *, section_n, sub_title, frame_phrase, tag,
                    active_stage, notes_id, frame_bar=False,
                    frame_size=20, illus=None, illus_caption=None,
                    passed=frozenset(), bar_muted=False):
    """v3.1 divider (#183 round 3, owner-мандат: «в заголовке каждого
    раздела должен быть мем или интересная картинка!»): 2-колоночная
    композиция — слева текстовый блок (номер раздела + подзаголовок +
    frame_phrase + tag), справа реальное фото/мем в Ocean rounded box,
    ЗАМЕТНОЕ (~35-38% площади слайда), единый паттерн на всех 6
    дивайдерах. pipeline_bar на всю ширину внизу, без изменений."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # --- Левая колонка: текст (ширина ~6.35", x=0.55..6.9) ---
    text_col_w = 6.35
    text_box(s, 0.55, 0.75, text_col_w, 1.9, f"Раздел {section_n}",
             size=92, bold=True, color=GOLD,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.55, 2.62, text_col_w, 1.05, sub_title,
             size=33, bold=True, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.08)
    text_box(s, 0.55, 3.80, text_col_w, 1.35, f"«{frame_phrase}»",
             size=frame_size, italic=True, color=MID, align=PP_ALIGN.LEFT,
             line_spacing=1.2)
    text_box(s, 0.55, 5.30, text_col_w, 0.45, tag,
             size=16, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    # --- Правая колонка: реальное фото/мем в Ocean rounded box
    #     (~5.4×4.55" ≈ 24.6 кв.дюйма ≈ 25% площади слайда — owner-мандат
    #     «≥25-30%, не декоративная малютка в углу») ---
    if illus is not None:
        img_x, img_y = 7.25, 0.85
        img_w, img_h = 5.55, 4.55
        ocean_box(s, img_x, img_y, img_w, img_h, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.5)
        pad = 0.16
        cap_h = 0.52 if illus_caption else 0
        _place_image_contain(s, illus, img_x + pad, img_y + pad,
                             img_w - 2 * pad, img_h - 2 * pad - cap_h)
        if illus_caption:
            text_box(s, img_x + pad, img_y + img_h - cap_h - 0.04,
                     img_w - 2 * pad, cap_h, illus_caption, size=9.5,
                     italic=True, color=LIGHT, align=PP_ALIGN.CENTER,
                     line_spacing=1.05)
    pipeline_bar(s, active_stage, passed=passed, muted_gold=bar_muted)
    if frame_bar:
        total_w = 12.3
        fx = (SLIDE_W_IN - total_w) / 2.0 - 0.14
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(fx), Inches(6.62 - 0.12),
                                 Inches(total_w + 0.28), Inches(0.48 + 0.24))
        shp.fill.background()
        shp.line.color.rgb = GOLD; shp.line.width = Pt(2.2)
        try:
            shp.adjustments[0] = 0.35
        except Exception:
            pass
        disable_shadow(shp)
    speaker_notes(s, load_notes(notes_id))
    return s


def _place_image_contain(slide, path, x, y, w, h):
    """Вписывает изображение в box (x,y,w,h) с сохранением пропорций
    (contain — не crop, не stretch), центрируя по обеим осям. Обходит
    #73-render-1 (add_picture с обоими w/h стрейчит непропорционально)."""
    if not Path(path).exists():
        return
    from PIL import Image as PILImage
    with PILImage.open(str(path)) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # шире бокса относительно — вписываем по ширине
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_y = y
        draw_x = x + (w - draw_w) / 2
    add_image(slide, path, x=draw_x, y=draw_y, w=draw_w)


# ============================================================
# Раздел 0
# ============================================================
def build_s01(p):
    """v3.1 мем-хук (#183 round 3, owner-мандат «вставь нормальные мемы и
    картинки из интернета»): реальный узнаваемый мем-шаблон «Well yes, but
    actually no» (imgflip, кадр из «Пираты! Банда неудачников», Aardman) —
    наш вопрос про T=0 в пустом верхнем поле шаблона, встроенная надпись
    мема снизу отвечает буквально на вопрос курса; hero ≥40% площади."""
    s = blank(p)
    # Hero: реальный мем-шаблон 1600×1218 (ratio 1.31), пустой верх ~21%
    # высоты — кладём туда свой вопрос текстовым слоем поверх картинки.
    hero_w = 7.6
    hero_h = hero_w * 1218 / 1600     # ≈ 5.79" (≈44% площади слайда)
    hx = (SLIDE_W_IN - hero_w) / 2
    hy = 0.55
    add_image(s, ASSETS / "web/well-yes-actually-no-template.jpg",
              x=hx, y=hy, w=hero_w)
    # Пустая белая полоса шаблона ≈ 0..21% высоты картинки
    blank_band_h = hero_h * 0.205
    text_box(s, hx + 0.35, hy + 0.10, hero_w - 0.7, blank_band_h - 0.15,
             "temperature=0 — значит, ответ всегда одинаковый?",
             size=27, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    # Встроенная в мем подпись «Well yes, but actually no» уже отвечает на
    # вопрос буквально — дополнительно подписываем расшифровку снизу.
    text_box(s, 0.55, hy + hero_h + 0.10, 12.23, 0.5,
             "Да, формально детерминирован. Но нет — на практике нет. "
             "Почему — сегодня.",
             size=16, italic=True, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Cover v2: большая «02» outline gold + title + subtitle новой арки +
    hero motif (4-стадийная конвейер-иконка справа)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    # «02» — outline gold, декоративная, справа
    outline_big_text(s, 9.15, 0.30, 4.1, 3.4, "02", size=230,
                     align=PP_ALIGN.CENTER)
    # Тег лекции
    text_box(s, 0.7, 1.05, 8.0, 0.5, "ЛЕКЦИЯ 2", size=18, bold=True,
             color=TEAL)
    filled_rect(s, 0.72, 1.58, 0.7, 0.05, TEAL)
    # Title
    text_box(s, 0.7, 2.0, 8.6, 2.3,
             "Как работают современные\nбольшие модели",
             size=44, bold=True, color=DEEP, line_spacing=1.08)
    # Subtitle — новая арка
    filled_rect(s, 0.7, 4.55, 0.05, 0.62, GOLD)
    text_box(s, 0.95, 4.55, 10.6, 0.75,
             "Конвейер инференса — и границы, которые меняют "
             "инженерные решения",
             size=20, italic=True, color=MID, line_spacing=1.22)
    # Hero motif: 4-стадийная конвейер-иконка (токен / вектор / внимание /
    # распределение)
    icons = [("[ ]", "токен"), ("0.21", "вектор"), ("⇄", "внимание"),
             ("%", "распределение")]
    pipe_y = 5.75
    cx = 0.95
    for i, (glyph, label) in enumerate(icons):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(pipe_y),
                                 Inches(0.78), Inches(0.78))
        shp.fill.solid()
        shp.fill.fore_color.rgb = MID if i != 2 else TEAL
        shp.line.fill.background(); disable_shadow(shp)
        tf = shp.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        r = para.add_run(); r.text = glyph
        r.font.name = FONT_MONO; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = WHITE
        text_box(s, cx - 0.45, pipe_y + 0.88, 1.7, 0.35, label,
                 size=11.5, color=LIGHT, align=PP_ALIGN.CENTER)
        if i < 3:
            right_arrow(s, cx + 0.92, pipe_y + 0.30, w=0.5, h=0.18,
                        fill=LIGHT)
        cx += 1.58
    speaker_notes(s, load_notes("s02"))


S02A_ROWS = [
    ("0", "Введение", "рамка и конвейер целиком", True),
    ("1", "Токенизация", "как модель видит ваш текст", False),
    ("2", "Эмбеддинги", "пространство смыслов и граница похожести", False),
    ("3", "Механизм внимания", "что важно сейчас: роли, кэш, длинный контекст", False),
    ("4", "Сэмплинг", "от распределения к токену: температура, детерминизм, невидимые токены", False),
    ("5", "Виды и размеры моделей", "на чём модели запускаются, мультимодальность, ландшафт 2026", False),
    ("6", "Финал", "сборка конвейера, итоги механизмов", False),
]


def build_s02a(p):
    """Карта лекции v3: 7 горизонтальных карточек-СТРОК (новый Раздел 5
    «Виды и размеры моделей»); без M-чипов, без минут; активный
    Раздел 0 — gold-обводка."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    slide_title(s, "Карта лекции — 7 разделов", size=28,
                align=PP_ALIGN.CENTER, y=0.35, h=0.62)
    row_w, row_h, gap = 12.1, 0.72, 0.10
    x0 = (SLIDE_W_IN - row_w) / 2
    y = 1.10
    for num, name, desc, active in S02A_ROWS:
        if active:
            ocean_box(s, x0, y, row_w, row_h, fill=WHITE, stroke=GOLD,
                      stroke_pt=2.5)
        else:
            ocean_box(s, x0, y, row_w, row_h, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.2)
        text_box(s, x0 + 0.25, y, 0.6, row_h, num, size=27, bold=True,
                 color=GOLD if active else LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        text_runs(s, x0 + 1.05, y, row_w - 1.4, row_h, [
            {"text": name, "size": 16, "bold": True, "color": DEEP},
            {"text": "  —  " + desc, "size": 13, "italic": True,
             "color": SLATE},
        ], anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + gap
    speaker_notes(s, load_notes("s02a"))


def build_s03(p):
    """v2.1: объект рассмотрения — слой «модель»; nested layers Лекции 1
    (Модель gold) + 2 строки Слой/Сегодня."""
    s = blank(p)
    slide_title(s, "Объект сегодня — слой «модель» из четырёх слоёв Лекции 1",
                size=26)
    # Nested layers слева, bottom-aligned (внешний = Приложение)
    base_x, bottom = 0.85, 6.75
    layers = [  # (label, w, h, fill, stroke, stroke_pt)
        ("Приложение", 5.5, 5.0, WHITE, LIGHT, 1.2),
        ("Агент", 4.6, 3.9, SURFACE, LIGHT, 1.2),
        ("Чат", 3.7, 2.8, TEAL_TINT, TEAL, 1.2),
        ("Модель", 2.8, 1.7, GOLD_TINT, GOLD, 2.2),
    ]
    for label, w, h, fill, stroke, sp in layers:
        x = base_x + (5.5 - w) / 2
        y = bottom - h
        ocean_box(s, x, y, w, h, fill=fill, stroke=stroke, stroke_pt=sp)
        text_box(s, x, y + 0.06, w, 0.4, label, size=14,
                 bold=(label == "Модель"),
                 color=DEEP if label == "Модель" else LIGHT,
                 align=PP_ALIGN.CENTER)
    # Справа — 2 строки
    ocean_box(s, 6.8, 2.0, 6.0, 1.6, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.3)
    text_runs(s, 7.0, 2.2, 5.6, 1.2, [
        {"text": "Слой «модель»: ", "size": 17, "bold": True, "color": MID},
        {"text": "stateless-инференс — на вход данные, на выход "
                 "предсказание, без памяти между вызовами.", "size": 15.5,
         "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 6.8, 4.0, 6.0, 1.6, GOLD_TINT, stroke=GOLD, stroke_pt=1.6,
                radius=True, radius_adj=0.10)
    text_runs(s, 7.0, 4.2, 5.6, 1.2, [
        {"text": "Сегодня: ", "size": 17, "bold": True, "color": DEEP},
        {"text": "разбираем, что происходит внутри этого инференса — и где "
                 "его устройство меняет инженерные решения.", "size": 15.5,
         "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """v3.0: цель лекции (gold на «важными деталями») + 7 промис-чипов
    по новой 7-разделной структуре (ряд 4 + ряд 3, без M-кодов)."""
    s = blank(p)
    slide_title(s, "Цель лекции", size=24, color=MID)
    ocean_box(s, 0.7, 1.25, 11.93, 2.05, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 1.05, 1.42, 11.25, 1.75, [
        {"text": "«Рассмотреть, как работает языковая модель, — и "
                 "разобраться с ", "size": 22, "bold": True, "color": DEEP},
        {"text": "важными деталями", "size": 22, "bold": True,
         "color": GOLD},
        {"text": ", которые меняют то, как вы строите промпты, агентов "
                 "и решения.»", "size": 22, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    promises = [
        "почему исправленный strawberry ничего не доказывает",
        "почему промпт с ролью реально меняет ответ",
        "что на деле умеет окно 1M",
        "почему T=0 не даёт одинаковых ответов",
        "сколько стоят невидимые токены рассуждения",
        "маленькая модель или гигант — по какому критерию выбирать",
        "чем заменить веру в бенчмарки",
    ]
    gap_x, gap_y = 0.24, 0.24
    card_h = 1.12
    # Ряд 1 — 4 карточки, ряд 2 — 3 карточки (центрированы)
    rows = [promises[:4], promises[4:]]
    y = 3.70
    for row_items in rows:
        n = len(row_items)
        card_w = (12.23 - gap_x * (n - 1)) / n if n == 4 else 3.93
        x = (SLIDE_W_IN - card_w * n - gap_x * (n - 1)) / 2
        for txt in row_items:
            ocean_box(s, x, y, card_w, card_h, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.2)
            filled_rect(s, x + 0.14, y + 0.17, 0.09, card_h - 0.34, GOLD)
            text_box(s, x + 0.38, y + 0.08, card_w - 0.55, card_h - 0.16,
                     txt, size=12.5, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.12)
            x += card_w + gap_x
        y += card_h + gap_y
    speaker_notes(s, load_notes("s04"))


def build_s04b(p):
    """Keystone: 7-стадийный конвейер инференса + петля авторегрессии
    (v2.1 #183: токены по одному, каждый добавляется ко входу) + пример +
    4 подкарточки + gold callout."""
    s = blank(p)
    slide_title(s, "Поток данных в LLM — туда и обратно", size=26)
    stages = [  # (name, example, caption)
        ("Текст", "«Привет»", "слова"),
        ("Токены", "[Прив][ет]", "id из словаря"),
        ("Векторы", "vec₁, vec₂", "числа"),
        ("LLM", "внимание", "инференс"),
        ("Распределение", "p(токен | контекст)", "вероятности"),
        ("Токен", "выбран", "выбор"),
        ("Текст", "ответ", "обратно в текст"),
    ]
    n = len(stages)
    arrow_w = 0.30
    total_w = 12.5
    cell_w = (total_w - arrow_w * (n - 1)) / n
    x0 = (SLIDE_W_IN - total_w) / 2
    y0, cell_h = 2.02, 1.35
    # Петля авторегрессии над конвейером: [Токен] (i=5) → назад к
    # [Токены] (i=1); подпись над линией
    x_from = x0 + 5 * (cell_w + arrow_w) + cell_w / 2
    x_to = x0 + 1 * (cell_w + arrow_w) + cell_w / 2
    loop_y = 1.74
    plain_line(s, x_from, y0, x_from, loop_y, color=TEAL, w_pt=2.2)
    plain_line(s, x_from, loop_y, x_to, loop_y, color=TEAL, w_pt=2.2)
    line_arrow(s, x_to, loop_y, x_to, y0 - 0.02, color=TEAL, w_pt=2.2)
    text_box(s, x_to + 0.25, loop_y - 0.36, x_from - x_to - 0.5, 0.32,
             "⟲ токены генерируются по одному; каждый добавляется ко входу",
             size=12, italic=True, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER)
    for i, (name, ex, cap) in enumerate(stages):
        x = x0 + i * (cell_w + arrow_w)
        is_llm = (name == "LLM")
        ocean_box(s, x, y0, cell_w, cell_h,
                  fill=GOLD_TINT if is_llm else SURFACE,
                  stroke=GOLD if is_llm else LIGHT,
                  stroke_pt=2.2 if is_llm else 1.4)
        nm_sz = 13 if len(name) < 12 else 10.5
        text_box(s, x - 0.06, y0 + 0.12, cell_w + 0.12, 0.4, name,
                 size=nm_sz, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        ex_sz = 10.5 if len(ex) < 14 else 8.5
        text_box(s, x - 0.10, y0 + 0.62, cell_w + 0.20, 0.6, ex,
                 size=ex_sz, color=MID, align=PP_ALIGN.CENTER,
                 font=FONT_MONO, line_spacing=1.1)
        text_box(s, x - 0.15, y0 + cell_h + 0.08, cell_w + 0.30, 0.5, cap,
                 size=10.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        if i < n - 1:
            right_arrow(s, x + cell_w + 0.03, y0 + cell_h / 2 - 0.08,
                        w=arrow_w - 0.06, h=0.16, fill=MID)
    # 4 подкарточки разделов
    subs = [("Раздел 1", "Текст → Токены"), ("Раздел 2", "Токены → Векторы"),
            ("Раздел 3", "LLM: внимание"), ("Раздел 4", "Распределение → Токен")]
    sub_w, sub_h, gap = 2.95, 0.95, 0.22
    sx0 = (SLIDE_W_IN - sub_w * 4 - gap * 3) / 2
    sy = 4.38
    for i, (nm, rng) in enumerate(subs):
        x = sx0 + i * (sub_w + gap)
        ocean_box(s, x, sy, sub_w, sub_h, fill=WHITE, stroke=TEAL,
                  stroke_pt=1.3)
        text_box(s, x, sy + 0.12, sub_w, 0.35, nm, size=13.5, bold=True,
                 color=TEAL, align=PP_ALIGN.CENTER)
        text_box(s, x, sy + 0.50, sub_w, 0.35, rng, size=12, color=DEEP,
                 align=PP_ALIGN.CENTER)
    gold_callout(s, 2.7, 5.78, 7.93, 0.75,
                 "Слова — только на границах; внутри — векторы.",
                 size=17, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s04b"))


# ============================================================
# Раздел 1 — Токенизация
# ============================================================
def build_s05a(p):
    section_divider(
        p, section_n=1, sub_title="Токенизация",
        frame_phrase="Как модель видит ваш текст",
        tag="3 разбора · 3 провала", active_stage=1, notes_id="s05a",
        passed={0},
        illus=ASSETS / "web/strawberry-openai-crop.jpg",
        illus_caption="Реальный диалог: ChatGPT о «strawberry» — "
                      "скриншот, OpenAI Community, 2024")


def token_chips_runs(pairs):
    """[(text, color)] -> run dicts моноширинно."""
    return [{"text": t, "size": 16, "bold": True, "color": c,
             "font": FONT_MONO} for t, c in pairs]


def build_s05(p):
    """3 примера разметки + gold callout + caption."""
    s = blank(p)
    slide_title(s, "Токен — id из словаря модели: не буква и не слово",
                size=26)
    rows = [
        ([("cat", DEEP)], [("[cat]", MID)], "1 токен"),
        ([("tokenization", DEEP)], [("[token]", MID), ("[ization]", TEAL)],
         "2 токена"),
        ([("клубника", DEEP)],
         [("[к]", MID), ("[луб]", TEAL), ("[ника]", LIGHT)],
         "3 токена (o200k_base)"),
    ]
    y = 1.75
    for src, toks, cnt in rows:
        ocean_box(s, 0.9, y, 11.5, 0.92, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.3)
        runs = token_chips_runs(src)
        runs.append({"text": "   →   ", "size": 16, "color": SLATE})
        runs += token_chips_runs(toks)
        runs.append({"text": "   →   ", "size": 16, "color": SLATE})
        runs.append({"text": cnt, "size": 17, "bold": True, "color": DEEP})
        text_runs(s, 1.25, y + 0.12, 10.9, 0.68, runs,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += 1.12
    gold_callout(s, 0.9, 5.35, 11.5, 0.75,
                 "«В среднем: 1 токен ≈ 4 символа EN ≈ 2 символа RU»",
                 size=17, align=PP_ALIGN.CENTER)
    text_box(s, 0.9, 6.35, 11.5, 0.5,
             "Словарь и модель — два разных артефакта: словарь строится до "
             "обучения модели, отдельным алгоритмом на своём корпусе.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """BPE: 2 колонки корпус -> словарь + gold callout + caption."""
    s = blank(p)
    slide_title(s, "BPE — компромисс между алфавитом и словарём", size=26)
    text_box(s, 0.55, 1.30, 12.3, 0.5,
             "Не все буквы (длинно) и не все слова (незнакомое выпадает) — "
             "частые подпоследовательности.",
             size=16, italic=True, color=MID)
    # Левая колонка — корпус (высота по контенту, v2.0.2 item 7)
    col_y, col_h = 2.05, 2.45
    ocean_box(s, 1.0, col_y, 4.7, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 1.3, col_y + 0.18, 4.1, 0.45, "Обучающий корпус",
             size=17, bold=True, color=MID)
    text_runs(s, 1.3, col_y + 0.75, 4.1, 2.0, [
        {"text": "low", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": "lower", "size": 16, "font": FONT_MONO, "color": DEEP,
         "newpara": True, "space_before_pt": 6},
        {"text": "newest", "size": 16, "font": FONT_MONO, "color": DEEP,
         "newpara": True, "space_before_pt": 6},
        {"text": "widest", "size": 16, "font": FONT_MONO, "color": DEEP,
         "newpara": True, "space_before_pt": 6},
    ])
    right_arrow(s, 6.05, col_y + col_h / 2 - 0.25, w=1.1, h=0.5, fill=MID)
    # Правая колонка — BPE-словарь
    ocean_box(s, 7.5, col_y, 4.7, col_h, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_box(s, 7.8, col_y + 0.18, 4.1, 0.45, "BPE-словарь",
             size=17, bold=True, color=TEAL)
    text_runs(s, 7.8, col_y + 0.75, 4.1, 2.0, [
        {"text": "low", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "er", "size": 16, "font": FONT_MONO, "color": TEAL},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "new", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "est", "size": 16, "font": FONT_MONO, "color": TEAL},
        {"text": " · ", "size": 16, "color": SLATE},
        {"text": "wid", "size": 16, "font": FONT_MONO, "color": DEEP},
        {"text": "+ одиночные символы, частые слоги, целые слова",
         "size": 12.5, "italic": True, "color": SLATE,
         "newpara": True, "space_before_pt": 12},
    ])
    gold_callout(s, 1.0, 4.95, 11.2, 0.80,
                 "Словарь строится один раз до обучения; на инференсе — "
                 "выборка готовых правил, не вычисление.", size=16,
                 align=PP_ALIGN.CENTER)
    text_box(s, 1.0, 6.05, 11.2, 0.5,
             "Разные вендоры режут один и тот же текст по-разному: Claude, "
             "GPT, Gemini — свои словари и правила.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s06"))


def build_s08(p):
    """v3.0 «Гонка патчей»: механизм слева + временная шкала патчей
    GPT-5.2 → GPT-5.5 → GPT-5.6 + StrawberryBench + callout."""
    s = blank(p)
    slide_title(s, "Гонка патчей: strawberry починили в апреле, "
                   "cranberry — только в июле", size=24)
    # Слева — механизм
    ocean_box(s, 0.55, 1.62, 5.0, 3.85, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.80, 4.4, 0.4, "Механизм — слепота к буквам",
             size=14.5, bold=True, color=MID)
    text_runs(s, 0.85, 2.28, 4.4, 0.95, [
        {"text": "strawberry", "size": 15, "font": FONT_MONO, "color": DEEP},
        {"text": "  →", "size": 15, "color": SLATE},
        {"text": "[st]", "size": 15, "font": FONT_MONO, "bold": True,
         "color": MID, "newpara": True, "space_before_pt": 6},
        {"text": "[raw]", "size": 15, "font": FONT_MONO, "bold": True,
         "color": TEAL},
        {"text": "[berry]", "size": 15, "font": FONT_MONO, "bold": True,
         "color": LIGHT},
    ])
    text_runs(s, 0.85, 3.32, 4.4, 0.4, [
        {"text": "Модель видит ", "size": 15, "color": DEEP},
        {"text": "3 токена", "size": 15, "bold": True, "color": DEEP},
        {"text": ", не 10 букв.", "size": 15, "color": DEEP},
    ])
    # Реальный скриншот: ChatGPT про strawberry (тот же кейс, что и
    # дивайдер s05a) — визуальное доказательство описанного механизма.
    # Используем tight-crop (1179×600, только ключевой обмен репликами).
    shot_w = 2.4
    shot_h = shot_w * 600 / 1179
    shot_x = 0.85 + (4.4 - shot_w) / 2
    shot_y = 3.75
    ocean_box(s, shot_x - 0.06, shot_y - 0.06, shot_w + 0.12, shot_h + 0.12,
              fill=RGBColor(0x11, 0x14, 0x18), stroke=TEAL, stroke_pt=1.2)
    add_image(s, ASSETS / "web/strawberry-openai-crop.jpg",
              x=shot_x, y=shot_y, w=shot_w)
    text_box(s, 0.85, shot_y + shot_h + 0.14, 4.4, 0.26,
             "скриншот: ChatGPT, форум OpenAI, 2024",
             size=9, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Справа — временная шкала гонки патчей (4 карточки вдоль оси)
    cards = [
        ("GPT-5.2 · дек 2025", "strawberry ✗ — «в strawberry две r»",
         False),
        ("GPT-5.5 · апр 2026",
         "strawberry ✓  /  cranberry ✗ — «две r» вместо трёх", False),
        ("GPT-5.6 · июль 2026",
         "cranberry ✓ — очередной вирусный кейс запатчен", True),
        ("StrawberryBench",
         "847 вопросов, 7 уровней сложности — системная проверка вместо "
         "вирусного вопроса", False),
    ]
    y = 1.55
    plain_line(s, 6.05, 1.8, 6.05, 5.15, color=LIGHT, w_pt=2.0)
    for title, body, is_gold in cards:
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.95),
                                 Inches(y + 0.40), Inches(0.2), Inches(0.2))
        shp.fill.solid()
        shp.fill.fore_color.rgb = GOLD if is_gold else LIGHT
        shp.line.fill.background(); disable_shadow(shp)
        if is_gold:
            filled_rect(s, 6.4, y, 6.4, 0.92, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.10)
        else:
            ocean_box(s, 6.4, y, 6.4, 0.92, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.2)
        text_box(s, 6.65, y + 0.07, 5.95, 0.36, title, size=13.5, bold=True,
                 color=DEEP if is_gold else MID)
        text_box(s, 6.65, y + 0.42, 5.95, 0.48, body, size=12, color=DEEP,
                 line_spacing=1.1)
        y += 1.0
    gold_callout(s, 0.55, 5.62, 12.25, 0.85,
                 "«Рваный интеллект» (jagged intelligence): уровень навыка "
                 "задаётся данными обучения, а не «общим умом» — модель "
                 "берёт олимпиадное золото и проваливает подсчёт букв.",
                 size=14.5, align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.60, 12.25, 0.4, [
        {"text": "Что делать: ", "size": 13.5, "bold": True, "color": MID},
        {"text": "тестируйте «cranberry своей предметки» — редкие случаи, "
                 "на которых никто не хайпил; вирусное прохождение ≠ навык.",
         "size": 13.5, "color": DEEP},
    ], align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """Числа и код: строка-мотивация (v2.1 #183) + 2 колонки + блок
    «Что делать» внизу."""
    s = blank(p)
    slide_title(s, "Токенизатор режет по частоте, а не по структуре",
                size=26, y=0.42, h=0.6)
    text_box(s, 0.55, 1.08, 12.3, 0.42,
             "Числа и код — самые частые «нетекстовые» входы: от их нарезки "
             "зависят арифметика модели и ваш бюджет токенов.",
             size=14.5, italic=True, color=MID)
    col_y, col_h = 1.62, 3.30
    # Колонка «Числа»
    ocean_box(s, 0.55, col_y, 6.0, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, col_y + 0.18, 5.4, 0.45, "Числа", size=18, bold=True,
             color=MID)
    text_runs(s, 0.85, col_y + 0.80, 5.45, 2.4, [
        {"text": "1000000", "size": 17, "font": FONT_MONO, "color": DEEP},
        {"text": " → ", "size": 17, "color": SLATE},
        {"text": "[100]", "size": 17, "font": FONT_MONO, "bold": True,
         "color": MID},
        {"text": "[000]", "size": 17, "font": FONT_MONO, "bold": True,
         "color": TEAL},
        {"text": "[0]", "size": 17, "font": FONT_MONO, "bold": True,
         "color": LIGHT},
        {"text": "Чанки слева направо — а математические разряды "
                 "считаются справа налево: границы не совпадают, сложение "
                 "в столбик ломается.",
         "size": 14, "color": DEEP, "newpara": True, "space_before_pt": 12},
        {"text": "Нарезка справа налево улучшает арифметику; "
                 "задача-специфичные схемы — ", "size": 14, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
        {"text": "до +33% точности", "size": 15, "bold": True,
         "color": GOLD},
        {"text": " к стандартной нарезке.", "size": 14, "color": DEEP},
    ], line_spacing=1.2)
    # Колонка «Код»
    ocean_box(s, 6.8, col_y, 6.0, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 7.1, col_y + 0.18, 5.4, 0.45, "Код", size=18, bold=True,
             color=TEAL)
    text_runs(s, 7.1, col_y + 0.80, 5.45, 2.4, [
        {"text": "GPT-2: ", "size": 15, "bold": True, "color": DEEP},
        {"text": "16 токенов", "size": 15.5, "bold": True, "color": MID},
        {"text": " на отступ 4-го уровня.", "size": 15, "color": DEEP},
        {"text": "GPT-4: ", "size": 15, "bold": True, "color": DEEP,
         "newpara": True, "space_before_pt": 12},
        {"text": "пробелы группами — словарь чинится, но не под все задачи "
                 "сразу.", "size": 15, "color": DEEP},
    ], line_spacing=1.22)
    # 3 приёма + заголовок группы (v2.0.2 item 9)
    text_box(s, 0.55, 5.10, 4.0, 0.32, "Что делать:", size=14, bold=True,
             color=MID)
    tips = ["Разделители разрядов («1 234 567»)",
            "Вычисления — в инструмент",
            "Единообразные отступы"]
    tip_w, gap = 3.95, 0.2
    x0 = (SLIDE_W_IN - tip_w * 3 - gap * 2) / 2
    for i, t in enumerate(tips):
        x = x0 + i * (tip_w + gap)
        ocean_box(s, x, 5.45, tip_w, 0.75, fill=WHITE, stroke=TEAL,
                  stroke_pt=1.3)
        text_box(s, x + 0.15, 5.51, tip_w - 0.3, 0.62, t, size=13,
                 bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    text_box(s, 0.55, 6.35, 12.25, 0.55,
             "Готовые чат-продукты уже сами уводят счёт во встроенные "
             "инструменты (code interpreter); звать инструмент самому — "
             "нестандартные случаи и свои приложения поверх API.",
             size=12, italic=True, color=MID, align=PP_ALIGN.CENTER,
             line_spacing=1.15)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    """Glitch-токены: story / механизм / факт + блок «Что на практике»
    (v2.1 #183: диагностика + санитизация) + footer-строка."""
    s = blank(p)
    slide_title(s, "Порядка 4% словаря — glitch-токены", size=26, y=0.42,
                h=0.6)
    col_y, col_h = 1.18, 3.42
    # Слева — story + механизм компактно (v3.0: анекдот свёрнут)
    ocean_box(s, 0.55, col_y, 3.95, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.8, col_y + 0.14, 3.45, 0.4, "SolidGoldMagikarp (2023)",
             size=13, bold=True, color=MID, font=FONT_MONO)
    text_box(s, 0.8, col_y + 0.56, 3.45, 0.65,
             "Юзернейм с Reddit в словаре GPT — модель не могла его "
             "повторить.", size=11.5, color=DEEP, line_spacing=1.12)
    # Реальный скриншот того же корпуса исследований (glitch-токен
    # «petertodd» — LessWrong, тот же автор/серия постов).
    shot2_w = 2.5
    shot2_h = shot2_w * 359 / 1157
    shot2_x = 0.8 + (3.45 - shot2_w) / 2
    shot2_y = col_y + 1.24
    ocean_box(s, shot2_x - 0.06, shot2_y - 0.06, shot2_w + 0.12,
              shot2_h + 0.12, fill=WHITE, stroke=TEAL, stroke_pt=1.1)
    add_image(s, ASSETS / "web/solidgoldmagikarp-1.png",
              x=shot2_x, y=shot2_y, w=shot2_w)
    text_runs(s, 0.8, shot2_y + shot2_h + 0.20, 3.45, 0.9, [
        {"text": "Механизм: ", "size": 11, "bold": True, "color": TEAL},
        {"text": "корпус словаря ≠ корпус модели → эмбеддинг токена "
                 "остаётся у случайной инициализации.", "size": 11,
         "color": DEEP},
    ], line_spacing=1.12)
    # Центр — «На что влияет» (v3.0: главный блок)
    ocean_box(s, 4.7, col_y, 3.95, col_h, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_box(s, 4.95, col_y + 0.16, 3.45, 0.45, "На что влияет",
             size=14.5, bold=True, color=TEAL)
    text_runs(s, 4.95, col_y + 0.62, 3.45, 2.7, [
        {"text": "• Сбои парсинга экзотических строк — редкие "
                 "идентификаторы, обфусцированный текст, необычный Unicode",
         "size": 12.5, "color": DEEP},
        {"text": "• Риски в проде — логи, user-generated ID, произвольный "
                 "пользовательский ввод", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
        {"text": "• Необъяснимые ответы «невпопад» без ошибки в коде",
         "size": 12.5, "color": DEEP, "newpara": True,
         "space_before_pt": 8},
    ], line_spacing=1.18)
    # Справа — факт GlitchMiner
    ocean_box(s, 8.85, col_y, 3.95, col_h, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 9.1, col_y + 0.16, 3.45, 0.45, "GlitchMiner (AAAI 2026)",
             size=14.5, bold=True, color=MID)
    text_runs(s, 9.1, col_y + 0.62, 3.45, 2.7, [
        {"text": "порядка 4% словаря", "size": 16, "bold": True,
         "color": GOLD},
        {"text": " по одной из оценок;", "size": 13.5, "color": DEEP},
        {"text": "воспроизводится в открытых семействах Llama, Qwen, "
                 "Gemma, Phi-3, Mistral.", "size": 13.5, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
    ], line_spacing=1.22)
    # Блок «Что на практике» — 2 карточки
    text_box(s, 0.55, 4.76, 5.0, 0.35, "Что на практике:", size=14,
             bold=True, color=MID)
    ocean_box(s, 0.55, 5.14, 6.0, 1.15, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 0.8, 5.24, 5.55, 0.95,
             "Необъяснимое поведение на экзотических строках (редкие "
             "идентификаторы, обфусцированный текст, необычный Unicode)? "
             "Гипотеза: glitch-токен. Проверка: замените строку "
             "плейсхолдером.",
             size=12.5, color=DEEP, line_spacing=1.14)
    ocean_box(s, 6.8, 5.14, 6.0, 1.15, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 7.05, 5.24, 5.55, 0.95,
             "В продуктах, принимающих произвольный ввод, — нормализация "
             "и санитизация входа до модели.",
             size=12.5, color=DEEP, line_spacing=1.14,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.55, 6.48, 12.25, 0.4,
             "Системная особенность конвейера, не баг версии — масштабом "
             "не чинится.",
             size=12.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """Стоимость языков: QuickChart bar слева + динамика и callout справа."""
    s = blank(p)
    slide_title(s, "Русский текст стоит ≈2× дороже английского", size=26)
    # Chart (980x560) — в ocean box
    box_x, box_y, box_w, box_h = 0.55, 1.7, 7.3, 4.55
    ocean_box(s, box_x, box_y, box_w, box_h, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    img_w = 6.9
    add_image(s, ASSETS / "charts/s11-tokens-per-char-v2.png",
              x=box_x + 0.2, y=box_y + 0.28, w=img_w)
    text_box(s, box_x + 0.2, box_y + 0.28 + img_w * 560 / 980 + 0.06,
             img_w, 0.35, "словари GPT-семейства",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Справа — динамика
    ocean_box(s, 8.15, 1.7, 4.65, 2.2, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.3)
    text_runs(s, 8.4, 1.9, 4.15, 1.8, [
        {"text": "Переход на o200k_base:", "size": 15, "bold": True,
         "color": MID},
        {"text": "примерно ", "size": 15, "color": DEEP, "newpara": True,
         "space_before_pt": 8},
        {"text": "−35%", "size": 17, "bold": True, "color": DEEP},
        {"text": " нелатинским языкам — разрыв сокращается, но не исчезает.",
         "size": 15, "color": DEEP},
    ], line_spacing=1.25)
    # Блок «Что делать» (v2.1 #183): калибровка на своём языке + когда
    # выгоден перевод на английский
    filled_rect(s, 8.15, 4.15, 4.65, 2.35, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.2, radius=True, radius_adj=0.08)
    text_runs(s, 8.38, 4.32, 4.2, 2.05, [
        {"text": "Что делать:", "size": 14, "bold": True, "color": MID},
        {"text": "•  Калибруйте лимиты в токенах на своём языке: фрагменты "
                 "для поиска, max_tokens, бюджет окна.", "size": 12.5,
         "color": DEEP, "newpara": True, "space_before_pt": 6},
        {"text": "•  Пакетная обработка больших объёмов — оцените перевод "
                 "на английский (≈2× дешевле); в интерактиве разница того "
                 "не стоит.", "size": 12.5, "color": DEEP, "newpara": True,
         "space_before_pt": 6},
    ], line_spacing=1.18)
    speaker_notes(s, load_notes("s11"))


# ============================================================
# Раздел 2 — Эмбеддинги
# ============================================================
def build_s12a(p):
    section_divider(
        p, section_n=2, sub_title="Эмбеддинги",
        frame_phrase="Пространство смыслов — и граница похожести",
        tag="4 разбора · 1 провал", active_stage=2, notes_id="s12a",
        passed={0, 1},
        illus=ASSETS / "web/word2vec-king-analogy-arrows.png",
        illus_caption="Классика word2vec: king − man + woman ≈ queen "
                      "(Jay Alammar, illustrated word2vec)")


def build_s12(p):
    """Эмбеддинг = выборка из входной таблицы."""
    s = blank(p)
    slide_title(s, "Каждому токену — вектор из входной таблицы модели",
                size=26)
    # Главная схема
    ocean_box(s, 0.55, 1.95, 8.0, 3.05, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    # [кот] чип
    chip(s, 0.95, 3.05, 1.1, 0.62, "[кот]", fill=MID, color=WHITE, size=17,
         font=FONT_MONO)
    right_arrow(s, 2.2, 3.22, w=0.55, h=0.28, fill=MID)
    # входная таблица — мини-грид 3x2
    tbl_x, tbl_y = 2.95, 2.50
    text_box(s, tbl_x - 0.15, tbl_y - 0.42, 2.6, 0.38, "входная таблица",
             size=13, bold=True, color=MID, align=PP_ALIGN.CENTER)
    for ri in range(3):
        for ci in range(2):
            fill = TEAL_TINT if ri == 1 else WHITE
            filled_rect(s, tbl_x + ci * 1.15, tbl_y + ri * 0.55, 1.1, 0.5,
                        fill, stroke=LIGHT, stroke_pt=1.0)
    text_box(s, tbl_x, tbl_y + 0.55, 1.1, 0.5, "[кот]", size=11,
             bold=True, color=DEEP, font=FONT_MONO, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, tbl_x + 1.15, tbl_y + 0.55, 1.1, 0.5, "→ вектор", size=10,
             color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    right_arrow(s, 5.5, 3.22, w=0.55, h=0.28, fill=MID)
    text_box(s, 6.15, 2.92, 2.3, 0.9, "[ 0.21, −0.45,\n0.88, …, 0.13 ]",
             size=14, bold=True, color=DEEP, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)
    text_box(s, 0.85, 4.45, 7.4, 0.45,
             "выучен на тренировке вместе с остальными весами; после "
             "обучения таблица фиксирована",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Mini-callout размерностей
    ocean_box(s, 8.85, 1.95, 3.95, 3.05, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 9.1, 2.15, 3.45, 0.4, "Размерности", size=15, bold=True,
             color=TEAL)
    text_runs(s, 9.1, 2.60, 3.45, 2.1, [
        {"text": "text-embedding-3-small — 1536", "size": 13, "color": DEEP,
         "font": FONT_MONO},
        {"text": "text-embedding-3-large — 3072", "size": 13, "color": DEEP,
         "font": FONT_MONO, "newpara": True, "space_before_pt": 6},
        {"text": "внутренние размерности флагманов не публикуются; "
                 "порядок — тысячи", "size": 12.5, "italic": True,
         "color": SLATE, "newpara": True, "space_before_pt": 8},
    ], line_spacing=1.2)
    gold_callout(s, 0.55, 5.55, 12.25, 0.62,
                 "«Геометрическая близость = смысловая близость»",
                 size=17, align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.28, 12.25, 0.5, [
        {"text": "Что делать: ", "size": 13, "bold": True, "color": TEAL},
        {"text": "опечатка или другой регистр — уже другой токен и другой "
                 "вектор; нормализуйте вход до эмбеддинга.", "size": 13,
         "color": DEEP},
    ], align=PP_ALIGN.CENTER, line_spacing=1.15)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """v2.1 (#183): эмбеддинги — не только внутренность инференса, но и
    инструмент поиска; три жизни от практики (3-я — gold, «ваш поиск/RAG»)
    + callout про переиндексацию."""
    s = blank(p)
    slide_title(s, "Эмбеддинги — не только внутренность инференса, "
                   "но и инструмент поиска", size=21.5, y=0.45, h=0.55)
    text_box(s, 0.55, 1.08, 12.3, 0.42,
             "Когда вы строите поиск или RAG, вы пользуетесь третьей жизнью "
             "термина «эмбеддинг» — отдельной embedding-моделью.",
             size=14.5, italic=True, color=MID)
    cards = [
        ("1", "Входная таблица", "внутри инференса",
         "Статическая: вектор [кот] один и тот же в любом предложении. "
         "Выборка по id, контекста ещё нет.", False),
        ("2", "Представление данных внутри модели", "внутри инференса",
         "Векторы, которые модель пересчитывает по ходу чтения контекста — "
         "после слоёв внимания. Именно они несут «понимание» модели.",
         False),
        ("3", "Векторы для поиска", "самостоятельный инструмент",
         "Вектор целого текста от отдельной embedding-модели — не "
         "внутренности вашего чат-LLM. Свой продукт, своё обучение, свои "
         "лидерборды.", True),
    ]
    card_w, col_h, gap = 3.95, 3.62, 0.2
    x0 = (SLIDE_W_IN - card_w * 3 - gap * 2) / 2
    y0 = 1.68
    for i, (num, title, sub, body, is_gold) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        if is_gold:
            filled_rect(s, x, y0, card_w, col_h, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.06)
        else:
            ocean_box(s, x, y0, card_w, col_h, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.4)
        text_box(s, x + 0.25, y0 + 0.14, 0.8, 0.65, num, size=32, bold=True,
                 color=GOLD if is_gold else LIGHT)
        text_box(s, x + 0.25, y0 + 0.82, card_w - 0.5, 0.8, title,
                 size=16, bold=True, color=DEEP, line_spacing=1.1)
        text_box(s, x + 0.25, y0 + 1.62, card_w - 0.5, 0.35, sub,
                 size=11.5, italic=True, color=MID if is_gold else SLATE)
        text_box(s, x + 0.25, y0 + 2.02, card_w - 0.5, 1.5, body,
                 size=13, color=DEEP, line_spacing=1.2)
        if is_gold:
            text_box(s, x + 0.25, y0 + col_h - 0.44, card_w - 0.5, 0.38,
                     "это и есть ваш поиск/RAG", size=11.5, italic=True,
                     bold=True, color=MID)
    gold_callout(s, 0.55, 5.72, 12.25, 0.85,
                 "Обновили чат-LLM → переиндексировать базу НЕ надо: индекс "
                 "живёт в координатах embedding-модели.",
                 size=16, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """v2.1 (#183): 2D scatter с подписанными осями-признаками (возврат
    иллюстрации размерности) + 3 факт-карточки."""
    s = blank(p)
    slide_title(s, "Близкие по смыслу токены лежат рядом — "
                   "в сотнях-тысячах измерений", size=24)
    # Слева — scatter в ocean box
    bx, by, bw, bh = 0.55, 1.65, 6.7, 5.0
    ocean_box(s, bx, by, bw, bh, fill=WHITE, stroke=LIGHT, stroke_pt=1.4)
    text_box(s, bx + 0.2, by + 0.12, 2.6, 0.35,
             "2D-проекция (PCA-стиль)", size=12.5, italic=True,
             color=SLATE)
    text_box(s, bx + 2.5, by + 0.12, bw - 2.7, 0.55,
             "каждая из 1536+ осей — выученный признак; здесь показаны две",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.RIGHT,
             line_spacing=1.05)
    # Оси-признаки (v2.1: возврат подписей осей)
    ax_x, ax_y = 1.02, 6.02  # начало координат (низ-лево)
    line_arrow(s, ax_x, ax_y, bx + bw - 0.25, ax_y, color=LIGHT, w_pt=1.8)
    line_arrow(s, ax_x, ax_y, ax_x, by + 0.62, color=LIGHT, w_pt=1.8)
    text_box(s, ax_x + 0.1, ax_y + 0.10, bw - 0.9, 0.32,
             "ось ≈ признак: тематика (веб-разработка ↔ кулинария)",
             size=10.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_box(s, ax_x + 0.12, by + 0.62, 3.6, 0.32,
             "ось ≈ признак: инфраструктура ↔ фронтенд",
             size=10.5, italic=True, color=LIGHT)
    # Точки
    def dot(x, y, fill, r=0.17):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                 Inches(r * 2), Inches(r * 2))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = WHITE; shp.line.width = Pt(1.2)
        disable_shadow(shp)
    # Пунктирные «облака» кластеров: SSL — верх-лево (инфраструктура),
    # React — низ-лево (фронтенд); борщ — справа (кулинария)
    for (ex, ey, ew, eh) in [(1.35, 2.55, 2.95, 1.5),
                             (1.55, 4.25, 3.05, 1.5)]:
        ell = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ex), Inches(ey),
                                 Inches(ew), Inches(eh))
        ell.fill.background()
        ell.line.color.rgb = LIGHT; ell.line.width = Pt(1.2)
        ell.line.dash_style = 4  # dash
        disable_shadow(ell)
    # SSL точки + подписи
    dot(1.80, 2.85, MID)
    text_box(s, 2.15, 2.72, 2.2, 0.55, "Как настроить SSL", size=11.5,
             bold=True, color=DEEP)
    dot(2.45, 3.42, MID)
    text_box(s, 2.90, 3.30, 2.3, 0.7, "Установка\nHTTPS-сертификата",
             size=11.5, bold=True, color=DEEP, line_spacing=1.05)
    # React точки + подписи
    dot(2.00, 4.55, TEAL)
    text_box(s, 2.35, 4.42, 2.3, 0.55, "Деплой React-компонента",
             size=11.5, bold=True, color=DEEP)
    dot(2.70, 5.12, TEAL)
    text_box(s, 3.14, 5.00, 2.2, 0.7, "Сборка\nReact-приложения",
             size=11.5, bold=True, color=DEEP, line_spacing=1.05)
    # Выброс — борщ (справа: кулинария)
    dot(5.65, 3.95, GOLD)
    text_box(s, 5.30, 4.35, 1.8, 0.4, "Рецепт борща", size=11.5,
             bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, 5.05, 4.62, 2.3, 0.4, "выброс — другая область",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # v3.0: иллюстративность пространства — «звёздная карта» (свой flat
    # рисунок, поддерживает scatter, не заменяет его)
    add_image(s, ASSETS / "illustrations/s14-space.png",
              x=5.25, y=2.35, w=1.70)
    # Справа — 3 факт-карточки
    facts = [
        ("Размерность", [
            {"text": "Публичные embedding-модели: ", "size": 13,
             "color": DEEP},
            {"text": "1536–3072", "size": 14.5, "bold": True, "color": GOLD},
            {"text": " измерения; внутренние у флагманов — порядка тысяч.",
             "size": 13, "color": DEEP}]),
        ("Обучение", [
            {"text": "Координаты не задаются вручную: похожие контексты "
                     "употребления → близкие векторы.", "size": 13,
             "color": DEEP}]),
        ("Проекция", [
            {"text": "Увидеть пространство можно только через PCA/t-SNE — "
                     "2D-картинка теряет часть структуры.", "size": 13,
             "color": DEEP}]),
    ]
    y = 1.65
    for title, runs in facts:
        ocean_box(s, 7.5, y, 5.3, 1.42, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.3)
        text_box(s, 7.75, y + 0.12, 4.8, 0.4, title, size=14.5, bold=True,
                 color=MID)
        text_runs(s, 7.75, y + 0.52, 4.8, 0.85, runs, line_spacing=1.18)
        y += 1.56
    gold_callout(s, 7.5, y + 0.08, 5.3, 1.0,
                 "Что делать: близость измерима расстоянием — фильтрация "
                 "и кластеризация без разметки и без LLM возможны прямо "
                 "на векторах, дёшево.", size=12.5, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s14"))


# ---- s15 heatmap (нативная PowerPoint-таблица, заливка Ocean-шкалой) ----
S15_LABELS_ROW = ["Как настроить SSL", "Установка HTTPS-сертификата",
                  "Деплой React-компонента", "Сборка React-приложения",
                  "Рецепт борща"]
S15_LABELS_COL = ["SSL", "HTTPS", "React-к.", "React-п.", "Борщ"]
S15_VALS = [
    [1.00, 0.85, 0.18, 0.20, 0.08],
    [0.85, 1.00, 0.22, 0.19, 0.07],
    [0.18, 0.22, 1.00, 0.78, 0.12],
    [0.20, 0.19, 0.78, 1.00, 0.10],
    [0.08, 0.07, 0.12, 0.10, 1.00],
]


def _ocean_scale(v):
    """0..1 -> RGB по шкале SURFACE -> LIGHT -> DEEP."""
    stops = [(0.0, (0xF4, 0xF7, 0xFA)), (0.5, (0x1C, 0x72, 0x93)),
             (1.0, (0x21, 0x29, 0x5C))]
    for (t0, c0), (t1, c1) in zip(stops, stops[1:]):
        if v <= t1:
            f = (v - t0) / (t1 - t0)
            return RGBColor(*(int(a + (b - a) * f) for a, b in zip(c0, c1)))
    return RGBColor(*stops[-1][1])


def build_s15(p):
    """Similarity-граница: heatmap 5x5 (таблица с Ocean-заливкой) +
    failure-карточка + callout."""
    s = blank(p)
    slide_title(s, "Высокое сходство — «об одном и том же», "
                   "не «с одинаковым смыслом»", size=24)
    # Таблица 6x6
    rows, cols = 6, 6
    tx, ty, tw, th = 0.55, 1.7, 7.15, 4.35
    gtbl = s.shapes.add_table(rows, cols, Inches(tx), Inches(ty),
                              Inches(tw), Inches(th))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(2.55)
    for ci in range(1, 6):
        tbl.columns[ci].width = Inches((tw - 2.55) / 5)
    tbl.rows[0].height = Inches(0.55)
    for ri in range(1, 6):
        tbl.rows[ri].height = Inches((th - 0.55) / 5)

    def cell_text(cell, txt, *, size=12, bold=False, color=DEEP,
                  align=PP_ALIGN.CENTER):
        cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        r = para.add_run(); r.text = txt
        r.font.name = FONT_BODY; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color

    # Угловая ячейка
    tbl.cell(0, 0).fill.solid()
    tbl.cell(0, 0).fill.fore_color.rgb = WHITE
    cell_text(tbl.cell(0, 0), "cosine similarity", size=10.5, bold=True,
              color=SLATE, align=PP_ALIGN.LEFT)
    # Заголовки колонок
    for ci, lab in enumerate(S15_LABELS_COL):
        c = tbl.cell(0, ci + 1)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=11.5, bold=True, color=MID)
    # Строки
    for ri, lab in enumerate(S15_LABELS_ROW):
        c = tbl.cell(ri + 1, 0)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=11.5, bold=True, color=MID,
                  align=PP_ALIGN.LEFT)
        for ci, v in enumerate(S15_VALS[ri]):
            cell = tbl.cell(ri + 1, ci + 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ocean_scale(v)
            txt_color = WHITE if v >= 0.45 else DEEP
            hot = (v in (0.85, 0.78))
            cell_text(cell, f"{v:.2f}".replace(".", ","), size=12,
                      bold=(v >= 0.7), color=txt_color)
    # Справа — failure-карточка gold
    filled_rect(s, 8.05, 1.7, 4.75, 2.6, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.2, radius=True, radius_adj=0.07)
    text_runs(s, 8.35, 1.95, 4.2, 2.1, [
        {"text": "«Как настроить SSL»", "size": 15.5, "bold": True,
         "color": DEEP},
        {"text": "  ↔", "size": 15.5, "bold": True, "color": MID},
        {"text": "«Как отключить SSL»", "size": 15.5, "bold": True,
         "color": DEEP, "newpara": True, "space_before_pt": 4},
        {"text": "Очень высокое сходство — противоположный практический "
                 "смысл.", "size": 14, "color": DEEP, "newpara": True,
         "space_before_pt": 10},
    ], line_spacing=1.2)
    # Мини-легенда шкалы
    text_box(s, 8.05, 4.5, 4.75, 0.35, "шкала: 0 — светлое · 1 — тёмное",
             size=11, italic=True, color=SLATE)
    gold_callout(s, 0.55, 6.3, 12.25, 0.8,
                 "Similarity — сигнал кандидатов; релевантность — отдельная "
                 "задача: реранкер, гибрид, фильтры.",
                 size=15.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s15"))


def build_s17(p):
    """Сборка раздела: путь туда-обратно + 2 карточки + callout + caption."""
    s = blank(p)
    slide_title(s, "Эмбеддинги — фундамент понимания: модель работает "
                   "с векторами, не строками", size=23)
    # Слева — вертикальная схема
    bx, by, bw, bh = 0.55, 1.6, 3.5, 5.35
    ocean_box(s, bx, by, bw, bh, fill=SURFACE, stroke=LIGHT, stroke_pt=1.4)
    steps = ["слова", "токены", "векторы", "LLM", "векторы", "токены",
             "слова"]
    step_h, ar_h = 0.52, 0.15
    total = len(steps) * step_h + (len(steps) - 1) * ar_h
    yy = by + (bh - total) / 2
    for i, st in enumerate(steps):
        is_llm = (st == "LLM")
        filled_rect(s, bx + 0.55, yy, bw - 1.1, step_h,
                    GOLD if is_llm else (WHITE if i % 2 == 0 else TEAL_TINT),
                    stroke=GOLD if is_llm else LIGHT,
                    stroke_pt=1.6 if is_llm else 1.0,
                    radius=True, radius_adj=0.35)
        text_box(s, bx + 0.55, yy + 0.05, bw - 1.1, step_h - 0.1, st,
                 size=14, bold=True, color=DEEP if not is_llm else DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            down_arrow(s, bx + bw / 2 - 0.09, yy + step_h + 0.01, w=0.18,
                       h=ar_h - 0.02, fill=MID)
        yy += step_h + ar_h
    # Справа — 2 карточки
    cards = [
        ("Перефразирования и синонимы",
         "«Как настроить SSL» и «Установка HTTPS-сертификата» — близкие "
         "векторы → модель отвечает одинаково; то же с «авто» и «машина»."),
        ("Межъязыковая близость",
         "«клубника» и strawberry — близкие векторы → ответ корректен "
         "независимо от языка запроса."),
    ]
    y = 1.6
    for title, body in cards:
        ocean_box(s, 4.4, y, 8.4, 1.55, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.3)
        text_box(s, 4.7, y + 0.13, 7.8, 0.4, title, size=15, bold=True,
                 color=MID)
        text_box(s, 4.7, y + 0.55, 7.8, 0.9, body, size=13.5, color=DEEP,
                 line_spacing=1.2)
        y += 1.75
    gold_callout(s, 4.4, 5.15, 8.4, 0.72,
                 "Семантическая близость на уровне предложений — основа "
                 "«понимания» переформулировок.", size=14.5)
    text_runs(s, 4.4, 5.97, 8.4, 0.75, [
        {"text": "Что делать: ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": "одна embedding-модель и индекс обслуживают поиск, "
                 "кластеризацию и RAG сразу — смена модели означает "
                 "переиндексацию всего хранилища; выбирайте как "
                 "инфраструктурное решение.", "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    text_box(s, 4.4, 6.85, 8.4, 0.45,
             "Выбор embedding-модели под задачу (MTEB, матрёшечные "
             "представления) — материал для самостоятельного чтения.",
             size=11.5, italic=True, color=LIGHT)
    speaker_notes(s, load_notes("s17"))


# ============================================================
# БАТЧ 2 — Раздел 3. Механизм внимания
# ============================================================
def build_s18a(p):
    section_divider(
        p, section_n=3, sub_title="Механизм внимания",
        frame_phrase="Как модель решает, на что опереться в контексте — и "
                     "что из этого следует для ролей, кэша и длинных окон",
        tag="4 разбора · 2 провала", active_stage=3, notes_id="s18a",
        passed={0, 1, 2},
        frame_size=18,
        illus=ASSETS / "web/attention-paper-title.png",
        illus_caption="Статья, давшая механизму имя: Vaswani et al., "
                      "«Attention Is All You Need», 2017 (arXiv:1706.03762)")


# v3.0: пример «Кот съел мышь, потому что ОН был голоден» — вес от «он»
# уходит к «Кот» (мужской род; раньше «она»→«мышь»).
S18_TOKENS = ["Кот", "съел", "мышь", "потому что", "он", "был", "голоден"]
S18_VALS = [
    [1.0, 0.3, 0.2, 0.1, 0.1, 0.1, 0.0],
    [0.4, 1.0, 0.5, 0.1, 0.1, 0.1, 0.1],
    [0.2, 0.4, 1.0, 0.1, 0.1, 0.0, 0.1],
    [0.1, 0.2, 0.2, 1.0, 0.3, 0.2, 0.2],
    [0.7, 0.1, 0.2, 0.2, 1.0, 0.3, 0.4],
    [0.1, 0.2, 0.2, 0.3, 0.4, 1.0, 0.5],
    [0.1, 0.2, 0.3, 0.3, 0.4, 0.5, 1.0],
]


def build_s18(p):
    """Attention-матрица 7×7 (нативная таблица, Ocean-шкала, «она»→«мышь»
    gold) + 3 свойства + callout."""
    s = blank(p)
    slide_title(s, "Внимание — это сверка каждого токена со всеми остальными",
                size=25, h=0.6)
    text_box(s, 0.55, 1.02, 12.3, 0.4,
             "Каждый токен «смотрит» на все остальные одновременно. "
             "На каждом шаге — N × N связей.",
             size=14, italic=True, color=MID)
    # Таблица 8×8
    rows, cols = 8, 8
    tx, ty, tw, th = 0.55, 1.55, 7.5, 3.95
    first_col = 1.30
    gtbl = s.shapes.add_table(rows, cols, Inches(tx), Inches(ty),
                              Inches(tw), Inches(th))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(first_col)
    for ci in range(1, 8):
        tbl.columns[ci].width = Inches((tw - first_col) / 7)
    tbl.rows[0].height = Inches(0.55)
    for ri in range(1, 8):
        tbl.rows[ri].height = Inches((th - 0.55) / 7)

    def cell_text(cell, txt, *, size=10.5, bold=False, color=DEEP,
                  align=PP_ALIGN.CENTER):
        cell.margin_left = Inches(0.02); cell.margin_right = Inches(0.02)
        cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        r = para.add_run(); r.text = txt
        r.font.name = FONT_BODY; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color

    c00 = tbl.cell(0, 0)
    c00.fill.solid(); c00.fill.fore_color.rgb = WHITE
    cell_text(c00, "вес", size=9.5, bold=True, color=SLATE)
    for ci, lab in enumerate(S18_TOKENS):
        c = tbl.cell(0, ci + 1)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=9.5, bold=True, color=MID)
    for ri, lab in enumerate(S18_TOKENS):
        c = tbl.cell(ri + 1, 0)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        cell_text(c, lab, size=10, bold=True, color=MID,
                  align=PP_ALIGN.LEFT)
        for ci, v in enumerate(S18_VALS[ri]):
            cell = tbl.cell(ri + 1, ci + 1)
            cell.fill.solid()
            is_gold = (ri == 4 and ci == 0)   # «он» → «Кот»
            is_future = (ci > ri)  # верхний треугольник — будущие токены
            if is_gold:
                cell.fill.fore_color.rgb = GOLD
                cell_text(cell, "0,7", size=11, bold=True, color=DEEP)
            elif is_future:
                # Засерён (v2.0.2 item 4): в декодере токен не видит будущих
                cell.fill.fore_color.rgb = SOFT_GREY
                cell_text(cell, f"{v:.1f}".replace(".", ","), size=9,
                          color=SLATE)
            else:
                cell.fill.fore_color.rgb = _ocean_scale(v)
                cell_text(cell, f"{v:.1f}".replace(".", ","), size=10,
                          bold=(v >= 0.7),
                          color=WHITE if v >= 0.45 else DEEP)
    text_runs(s, 0.55, 5.62, 7.5, 0.8, [
        {"text": "По строке «он»: ", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": "наибольший вес — на «Кот» (мужской род). "
                 "Статистическая связь, выученная на корпусе.",
         "size": 12.5, "color": DEEP},
        {"text": "В декодере токен видит только предыдущие — показана "
                 "полная сверка для наглядности.", "size": 11.5,
         "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 3},
    ], line_spacing=1.12)
    # Справа — 3 свойства
    props = [
        ("Размерность", [
            {"text": "N × N, где N — длина контекста. Удвоение контекста — ",
             "size": 12.5, "color": DEEP},
            {"text": "учетверение вычислений внимания", "size": 12.5,
             "bold": True, "color": DEEP},
            {"text": ".", "size": 12.5, "color": DEEP}]),
        ("На каждом шаге", [
            {"text": "Распределение весов пересчитывается заново на каждом "
                     "шаге генерации.", "size": 12.5, "color": DEEP}]),
        ("Многоголовость", [
            {"text": "В каждом слое — десятки параллельных «голов» (типично "
                     "32–128); каждая ловит свой тип связей: грамматика, "
                     "семантика, дальние зависимости.", "size": 12.5,
             "color": DEEP}]),
    ]
    y = 1.55
    for title, runs in props:
        ocean_box(s, 8.3, y, 4.5, 1.48, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.3)
        text_box(s, 8.55, y + 0.10, 4.0, 0.38, title, size=14, bold=True,
                 color=MID)
        text_runs(s, 8.55, y + 0.50, 4.0, 0.9, runs, line_spacing=1.15)
        y += 1.62
    gold_callout(s, 0.55, 6.45, 12.25, 0.70,
                 "На что влияют веса: они определяют, чьи Value попадут в "
                 "представление текущего токена — и напрямую формируют "
                 "следующее предсказание.",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    """v3.0: Q/K/V на НАШЕМ предложении («он был голоден») — worked-пример
    слева, chart весов (лидер «Кот» gold) справа + 3 факта + фонарик."""
    s = blank(p)
    slide_title(s, "Внимание — распределение весов на весь контекст: "
                   "три проекции Query / Key / Value", size=20, h=0.65,
                y=0.35)
    # Главный тезис — сразу под заголовком (gold)
    gold_callout(s, 0.55, 1.06, 12.25, 0.6,
                 "Q — про текущий шаг. K и V — про уже обработанный контекст.",
                 size=15.5, align=PP_ALIGN.CENTER)
    # Строка Q/K/V — три плашки
    qkv = [("Query", "«что я сейчас ищу»"),
           ("Key", "«что я предлагаю»"),
           ("Value", "«что я отдам, если меня выбрали»")]
    x = 0.55
    for term, phrase in qkv:
        ocean_box(s, x, 1.84, 3.95, 0.58, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, x + 0.18, 1.84, 3.65, 0.58, [
            {"text": term, "size": 13.5, "bold": True, "color": MID,
             "font": FONT_MONO},
            {"text": " — " + phrase, "size": 12, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        x += 4.15
    # Слева — worked-пример на нашем предложении
    ocean_box(s, 0.55, 2.62, 6.35, 3.05, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.3)
    text_runs(s, 0.85, 2.76, 5.8, 0.4, [
        {"text": "«Кот съел мышь, потому что ", "size": 14.5, "color": DEEP},
        {"text": "он", "size": 14.5, "bold": True, "color": MID},
        {"text": " был голоден»", "size": 14.5, "color": DEEP},
    ])
    text_runs(s, 0.85, 3.30, 5.8, 1.6, [
        {"text": "Q(«он»)", "size": 13, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": " = «ищу: кто мог быть голоден»", "size": 12.5,
         "color": DEEP},
        {"text": "K(«Кот»)", "size": 13, "bold": True, "color": TEAL,
         "font": FONT_MONO, "newpara": True, "space_before_pt": 7},
        {"text": " = «я — одушевлённый субъект»", "size": 12.5,
         "color": DEEP},
        {"text": "V(«Кот»)", "size": 13, "bold": True, "color": TEAL,
         "font": FONT_MONO, "newpara": True, "space_before_pt": 7},
        {"text": " = содержимое, которое вольётся в представление «он»",
         "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    text_box(s, 0.85, 4.95, 5.8, 0.65,
             "Совпадение Q(«он») и K(«Кот») даёт высокий вес → V(«Кот») "
             "определяет обновлённое представление токена «он».",
             size=11.5, italic=True, color=MID, line_spacing=1.15)
    # Справа — chart (7 токенов, та же нарезка предложения, что в s18)
    ocean_box(s, 7.05, 2.62, 5.75, 3.05, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 7.3, 2.74, 5.25, 0.35,
             "Распределение весов — сумма = 1, лидер «Кот»",
             size=12.5, bold=True, color=MID)
    add_image(s, ASSETS / "charts/s19-attention-weights.png",
              x=7.3, y=3.18, w=5.25)
    # 3 нумерованных факта — горизонтальный ряд
    triples = [
        ("1", "На вход — все токены контекста."),
        ("2", "На выходе — распределение весов, сумма = 1."),
        ("3", "Пересчитывается на каждом шаге генерации."),
    ]
    x = 0.55
    for num, txt in triples:
        ocean_box(s, x, 5.85, 3.95, 0.62, fill=WHITE, stroke=TEAL,
                  stroke_pt=1.2)
        text_runs(s, x + 0.18, 5.85, 3.65, 0.62, [
            {"text": num + ".  ", "size": 13, "bold": True, "color": TEAL},
            {"text": txt, "size": 11.5, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
        x += 4.15
    # Метафора фонарика — одна строка-подпись
    text_box(s, 0.55, 6.62, 12.3, 0.45,
             "Метафора: фонарик в тёмной комнате — луч направлен на "
             "релевантные токены, яркость света = вес внимания.",
             size=12.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    """v3.0: worked example «он»→«Кот» (веер стрелок) + контраст с ролью и
    без («Объясни GIL») + research-оговорка (Zheng et al. — в notes) +
    gold callout. Спуфинг убран (тема → Лекция 3)."""
    s = blank(p)
    slide_title(s, "Роль работает через вес во внимании — но не повышает "
                   "фактическое качество", size=20, h=0.65, y=0.32)
    # Верхний box: рабочий пример
    ocean_box(s, 0.55, 1.08, 12.25, 2.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.20, 5.5, 0.4, "Рабочий пример: куда смотрит «он»",
             size=14, bold=True, color=MID)
    # Токены предложения — отдельные боксы для стрелок
    seg = [("«Кот", 2.35, 0.85, True),
           ("съел мышь,", 3.25, 1.60, False),
           ("потому что", 4.90, 1.70, False),
           ("он", 6.70, 0.55, True),
           ("был", 7.30, 0.80, False),
           ("голоден»", 8.15, 1.55, False)]
    for txt, x, w, bold in seg:
        text_box(s, x, 1.66, w, 0.45, txt, size=19, bold=bold,
                 color=DEEP if not bold else MID)
    # Стрелки от «он» (x≈6.85) к целям: «Кот» (gold), «был», «голоден»
    line_arrow(s, 6.85, 2.13, 2.65, 2.50, color=GOLD, w_pt=4.0)
    line_arrow(s, 7.00, 2.13, 7.55, 2.50, color=MID, w_pt=2.2)
    line_arrow(s, 7.10, 2.13, 8.85, 2.50, color=LIGHT, w_pt=1.2)
    text_box(s, 2.65, 2.56, 2.4, 0.32, "главный вес", size=11.5, bold=True,
             color=GOLD)
    text_box(s, 0.85, 2.90, 6.8, 0.42,
             "Упрощение: агрегат сотен связей в десятках слоёв — модель "
             "воспроизводит корреляции употребления, не грамматический "
             "разбор.", size=10.5, italic=True, color=SLATE,
             line_spacing=1.1)
    text_runs(s, 7.85, 2.68, 4.7, 0.62, [
        {"text": "Подумайте 30 секунд: ", "size": 11.5, "bold": True,
         "color": DEEP},
        {"text": "куда уйдёт вес от «он» в «Программа упала, потому что он "
                 "забыл обработать null»?", "size": 11.5, "color": DEEP},
    ], line_spacing=1.1)
    # Контраст: без роли / с ролью
    ocean_box(s, 0.55, 3.48, 5.7, 1.5, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 0.85, 3.60, 5.1, 0.38, "Без роли", size=15, bold=True,
             color=MID)
    text_runs(s, 0.85, 4.00, 5.15, 0.9, [
        {"text": "«Объясни GIL»", "size": 12.5, "bold": True, "color": DEEP},
        {"text": " → нейтральный, обобщённый ответ.", "size": 12.5,
         "color": DEEP},
    ], line_spacing=1.18)
    right_arrow(s, 6.38, 4.05, w=0.6, h=0.34, fill=GOLD)
    ocean_box(s, 7.1, 3.48, 5.7, 1.5, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 7.4, 3.60, 5.1, 0.38, "С ролью", size=15, bold=True,
             color=TEAL)
    text_runs(s, 7.4, 4.00, 5.15, 0.95, [
        {"text": "«Ты — опытный Python-разработчик", "size": 12.5,
         "bold": True, "color": TEAL},
        {"text": ". Объясни GIL» — токены роли получают вес → смещают "
                 "распределение следующих токенов: конкретнее, в экспертном "
                 "регистре.", "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    # Research-оговорка (fact-check fix, 2026-09-06): Zheng et al. измеряли
    # ТОЛЬКО фактическую точность (2410 вопросов, 162 роли) — про
    # тон/стиль в их статье вывода нет. Разделяем на 2 отдельные строки:
    # (1) точная цитата с цифрами, (2) наблюдение курса про тон/стиль —
    # без привязки к Zheng et al.
    ocean_box(s, 0.55, 5.05, 12.25, 0.82, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 0.85, 5.11, 11.7, 0.36,
             "Zheng et al. (2024, EMNLP Findings; 2410 вопросов, 162 роли): "
             "персона/роль в промпте не повышает точность факта — эффект "
             "конкретной роли непредсказуем.",
             size=12.5, bold=True, color=DEEP, line_spacing=1.1)
    text_box(s, 0.85, 5.50, 11.7, 0.32,
             "Отдельно от этого исследования — по наблюдениям курса: роль "
             "ощутимо меняет тон, стиль и отбор содержания ответа.",
             size=11.5, italic=True, color=TEAL, line_spacing=1.05)
    gold_callout(s, 0.55, 6.0, 12.25, 0.95,
                 "Роль — инструмент управления стилем и фокусом, а не "
                 "«усилитель ума». Нужен ответ по вашим данным — дайте "
                 "данные, а не третье прилагательное к слову «эксперт».",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s20"))


def build_s21(p):
    """KV-cache: схема K/V в кэше + Q нового токена; фазы prefill/decode;
    gold callout про длинный чат."""
    s = blank(p)
    slide_title(s, "K и V кешируются — заново считается только Q", size=26,
                h=0.6)
    text_runs(s, 0.55, 1.08, 12.3, 0.55, [
        {"text": "KV-cache: ", "size": 14.5, "bold": True, "color": MID},
        {"text": "Key/Value обработанных токенов сохраняются в памяти "
                 "ускорителя. На каждом шаге вычисляется только Q нового "
                 "токена — против сохранённых K/V.", "size": 14.5,
         "color": DEEP},
    ], line_spacing=1.15)
    # Схема
    ocean_box(s, 0.55, 1.75, 12.25, 1.85, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    xx = 1.0
    for i in range(4):
        lab = f"токен {i+1}" if i < 3 else "…"
        chip(s, xx, 2.0, 1.15, 0.42, lab, fill=MID, color=WHITE, size=12)
        for j, kv in enumerate(("K", "V")):
            filled_rect(s, xx + 0.06 + j * 0.56, 2.52, 0.5, 0.42, TEAL_TINT,
                        stroke=TEAL, stroke_pt=1.2, radius=True,
                        radius_adj=0.25)
            text_box(s, xx + 0.06 + j * 0.56, 2.56, 0.5, 0.34, kv, size=13,
                     bold=True, color=TEAL, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
        xx += 1.45
    text_box(s, 1.0, 3.06, 5.6, 0.35,
             "уже посчитаны — лежат в кэше, не пересчитываются",
             size=11, italic=True, color=TEAL)
    # Разделитель и новый токен
    plain_line(s, 7.15, 1.95, 7.15, 3.35, color=LIGHT, w_pt=1.2, dash=4)
    chip(s, 9.7, 2.0, 1.55, 0.42, "новый токен", fill=DEEP, color=WHITE,
         size=12)
    filled_rect(s, 10.2, 2.52, 0.5, 0.42, GOLD, stroke=None, radius=True,
                radius_adj=0.25)
    text_box(s, 10.2, 2.56, 0.5, 0.34, "Q", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    left_arrow(s, 7.45, 2.6, w=2.6, h=0.26, fill=GOLD)
    text_box(s, 7.45, 2.95, 2.7, 0.35, "сверка со всеми K/V",
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)
    # Фазы
    ocean_box(s, 0.55, 3.72, 5.9, 1.95, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 0.85, 3.84, 5.3, 0.4, "Фаза 1 — prefill (обработка промпта)",
             size=14, bold=True, color=MID)
    text_runs(s, 0.85, 4.26, 5.35, 1.35, [
        {"text": "• Все токены входа известны сразу → K/V считаются ",
         "size": 12.5, "color": DEEP},
        {"text": "параллельно", "size": 12.5, "bold": True, "color": DEEP},
        {"text": "• Упирается в ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "вычислительную мощность", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": "• Определяет ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "паузу до первого символа ответа (TTFT)", "size": 12.5,
         "bold": True, "color": DEEP},
    ], line_spacing=1.14)
    right_arrow(s, 6.55, 4.55, w=0.45, h=0.28, fill=MID)
    ocean_box(s, 7.1, 3.72, 5.7, 1.95, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 7.4, 3.84, 5.1, 0.4, "Фаза 2 — decode (генерация ответа)",
             size=14, bold=True, color=TEAL)
    text_runs(s, 7.4, 4.26, 5.15, 1.35, [
        {"text": "• Строго ", "size": 12.5, "color": DEEP},
        {"text": "последовательно", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": ", токен за токеном", "size": 12.5, "color": DEEP},
        {"text": "• Каждый шаг читает ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "весь накопленный кэш", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": " из памяти", "size": 12.5, "color": DEEP},
        {"text": "• Упирается в ", "size": 12.5, "color": DEEP,
         "newpara": True, "space_before_pt": 5},
        {"text": "пропускную способность памяти", "size": 12.5,
         "bold": True, "color": DEEP},
        {"text": " → скорость «печати»", "size": 12.5, "color": DEEP},
    ], line_spacing=1.14)
    # Вендоры: где кэш автоматический, где явный
    ocean_box(s, 0.55, 5.80, 12.25, 0.60, fill=SURFACE, stroke=TEAL,
              stroke_pt=1.2)
    text_runs(s, 0.85, 5.80, 11.7, 0.60, [
        {"text": "Кэш у провайдеров:  ", "size": 12.5, "bold": True,
         "color": TEAL},
        {"text": "OpenAI — авто  ·  Gemini — implicit (авто)  ·  "
                 "DeepSeek — авто-дисковый  ·  Anthropic — явный "
                 "(cache_control)", "size": 12.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.52, 12.25, 0.75,
                 "KV-cache делает повторную подачу истории дешёвой и "
                 "быстрой; «тормозит и дорожает» — когда кэш НЕ срабатывает "
                 "(длинный чат без кэша, нестабильный префикс).",
                 size=13.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """v3.0 prompt caching: схема ТРЁХ запросов (запись → hit → miss) +
    кейс-бары слева; exact-prefix стек справа; gold callout. Mini-poll
    снят (интерактива нет)."""
    s = blank(p)
    slide_title(s, "Кэш промптов — ставка на повторное использование, "
                   "а не скидка", size=25, h=0.6)
    # Слева — схема трёх запросов
    ocean_box(s, 0.55, 1.3, 6.35, 4.35, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.42, 5.8, 0.38, "Три запроса подряд:",
             size=14.5, bold=True, color=MID)
    reqs = [
        ("Запрос 1", "префикс пишется в кэш — запись 1.25–2× ставки",
         TEAL_TINT, TEAL, False),
        ("Запрос 2", "тот же префикс → cache hit — чтение 0.1× "
         "(новейшие — 0.025×)", GOLD_TINT, GOLD, True),
        ("Запрос 3", "строка добавлена в начало → префикс не совпал → "
         "мимо кэша, полная цена", SOFT_GREY, SLATE, False),
    ]
    yy = 1.86
    for name, desc, fill, stroke, is_gold in reqs:
        filled_rect(s, 0.85, yy, 5.75, 0.62, fill, stroke=stroke,
                    stroke_pt=1.6 if is_gold else 1.2, radius=True,
                    radius_adj=0.14)
        text_runs(s, 1.02, yy + 0.03, 5.45, 0.56, [
            {"text": name + ":  ", "size": 12, "bold": True,
             "color": DEEP},
            {"text": desc, "size": 11.5, "color": DEEP,
             "bold": is_gold},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
        yy += 0.72
    # Кейс — 2 бара
    text_box(s, 0.85, 4.12, 5.8, 0.35,
             "Кейс: 50 000 анализов документов в месяц", size=13, bold=True,
             color=DEEP)
    filled_rect(s, 0.85, 4.50, 4.6, 0.32, MID)
    text_box(s, 5.55, 4.50, 1.3, 0.32, "$45 000", size=12, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.98, 4.52, 3.6, 0.28, "без кэша", size=11, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 0.85, 4.92, 0.82, 0.32, GOLD)
    text_runs(s, 1.8, 4.92, 4.9, 0.32, [
        {"text": "$8 000 с кэшем · ", "size": 12, "bold": True,
         "color": DEEP},
        {"text": "−82%", "size": 13.5, "bold": True, "color": GOLD},
    ], anchor=MSO_ANCHOR.MIDDLE)
    # Справа — exact prefix (иллюстрирует запрос 3)
    ocean_box(s, 7.15, 1.3, 5.65, 4.35, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 7.4, 1.45, 5.15, 0.75, [
        {"text": "Условие: точное совпадение префикса. ", "size": 13,
         "bold": True, "color": DEEP},
        {"text": "Один изменившийся токен инвалидирует кэш для всего, "
                 "что после него.", "size": 12.5, "color": DEEP},
    ], line_spacing=1.15)
    stack = [
        ("системный промпт", TEAL_TINT, TEAL, "✓ кэш"),
        ("инструкции", TEAL_TINT, TEAL, "✓ кэш"),
        ("документы", TEAL_TINT, TEAL, "✓ кэш"),
        ("«Сегодня 2026-09-05 14:23» — динамический", GOLD_TINT, GOLD, "!"),
        ("вопрос", SOFT_GREY, SLATE, "× мимо кэша"),
    ]
    yy = 2.3
    for txt, fill, stroke, mark in stack:
        filled_rect(s, 7.4, yy, 4.0, 0.42, fill, stroke=stroke,
                    stroke_pt=1.4, radius=True, radius_adj=0.2)
        text_box(s, 7.55, yy + 0.02, 3.75, 0.38, txt,
                 size=10.5 if len(txt) > 20 else 12, bold=(mark == "!"),
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 11.5, yy + 0.02, 1.2, 0.38, mark, size=11, bold=True,
                 color=stroke if mark != "!" else GOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.52
    text_box(s, 7.4, 4.95, 5.15, 0.6,
             "Динамический элемент в начале (timestamp) рушит кэш всех "
             "блоков после него — это и есть «запрос 3».",
             size=11.5, italic=True, color=SLATE, line_spacing=1.12)
    gold_callout(s, 0.55, 5.95, 12.25, 0.8,
                 "Правило компоновки: стабильное — в начало (промпт, "
                 "инструкции, примеры, документы), переменное — в конец.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    """Гонка окон: лог-бары 4 модели + 2 полюса + строка про позиционное
    кодирование + gold callout."""
    s = blank(p)
    slide_title(s, "Фронтир-стандарт 2026 — окно до 1 млн токенов. "
                   "Но заявленное ≠ полезное", size=24, h=0.9, y=0.35)
    # Лог-бары
    ocean_box(s, 0.55, 1.5, 7.65, 3.55, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    bars = [
        ("GPT-3.5 (2022)", "4 тыс.", 4_000, MID, ""),
        ("Claude 3.5 (2024)", "200 тыс.", 200_000, MID, ""),
        ("Фронтир-стандарт (2026)", "до 1 млн", 1_000_000, GOLD,
         "Fable 5, GPT-5.6, Gemini 3.1 Pro и др."),
        ("Gemini 3.5 Pro (2026)", "2 млн", 2_000_000, MID,
         "единичная модель, не стандарт"),
    ]
    import math
    yy = 1.68
    for name, val, v, color, note in bars:
        w = (math.log10(v) - 3.0) / 3.4 * 4.5
        text_runs(s, 0.9, yy, 6.9, 0.3, [
            {"text": name + " — ", "size": 12.5, "bold": True,
             "color": DEEP},
            {"text": val, "size": 12.5, "bold": True,
             "color": GOLD if color == GOLD else MID},
            {"text": ("  ·  " + note) if note else "", "size": 11,
             "italic": True, "color": SLATE},
        ])
        filled_rect(s, 0.9, yy + 0.33, max(w, 0.3), 0.26, color,
                    radius=True, radius_adj=0.35)
        yy += 0.78
    text_box(s, 0.9, 4.73, 6.9, 0.3, "ширина — логарифмическая шкала",
             size=10.5, italic=True, color=SLATE)
    # Полюса
    card = filled_rect(s, 8.45, 1.5, 4.35, 1.68, SOFT_GREY, stroke=LIGHT,
                       stroke_pt=1.4, radius=True, radius_adj=0.08)
    card.line.dash_style = 4
    text_box(s, 8.7, 1.65, 3.85, 0.4, "Llama 4 Scout: «10 млн»", size=14,
             bold=True, color=DEEP)
    text_box(s, 8.7, 2.1, 3.85, 1.0,
             "заявлено; ни один опубликованный бенчмарк не подтверждает "
             "качество вблизи предела", size=12, color=SLATE,
             line_spacing=1.18)
    ocean_box(s, 8.45, 3.37, 4.35, 1.68, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 8.7, 3.52, 3.85, 0.4, "YandexGPT 5 Pro: 32 тыс.", size=14,
             bold=True, color=DEEP)
    text_box(s, 8.7, 3.97, 3.85, 1.0,
             "на полтора-два порядка меньше флагманов — для длинных "
             "документов это определяющее ограничение", size=12, color=DEEP,
             line_spacing=1.18)
    text_box(s, 0.55, 5.2, 12.25, 0.55,
             "Просто «растянуть» окно нельзя: позиция токена закодирована "
             "геометрией, обученной на конкретных длинах, — расширение "
             "(RoPE / YaRN) — отдельная инженерная работа.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER,
             line_spacing=1.15)
    gold_callout(s, 0.55, 5.9, 12.25, 0.68,
                 "Платите за то, что кладёте в окно, а не за то, что окно "
                 "вмещает: 900 тыс. токенов входа по $10/млн ≈ $9 за один "
                 "вызов.", size=13.5, align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.63, 12.25, 0.55, [
        {"text": "Что делать: ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": "выбирайте модель по эффективному окну задачи (бенчмарки "
                 "без лексических подсказок), не по маркетинговому "
                 "максимуму.", "size": 12.5, "color": DEEP},
    ], align=PP_ALIGN.CENTER, line_spacing=1.12)
    speaker_notes(s, load_notes("s23"))


def build_s25(p):
    """v3.0 редизайн в ДВА ЯРУСА (верх/низ): верх — «забывание победили»
    (U-кривая 2023 → плоская линия 2026, свой chart); низ — «работу со
    смыслом — нет» (NoLiMa-панель с явной подписью). + формула-callout."""
    s = blank(p)
    slide_title(s, "Поиск дословной вставки решён. Понимание длинного "
                   "контекста — нет", size=24, h=0.62, y=0.32)
    # ── ВЕРХНИЙ ЯРУС: забывание победили
    ocean_box(s, 0.55, 1.02, 12.25, 2.42, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.85, 1.12, 11.7, 0.4, [
        {"text": "Забывание победили ", "size": 14.5, "bold": True,
         "color": MID},
        {"text": "✓", "size": 16, "bold": True, "color": TEAL},
        {"text": "  — найти дословную вставку (needle-in-a-haystack)",
         "size": 13, "color": DEEP},
    ])
    add_image(s, ASSETS / "charts/s25-ucurve.png", x=0.85, y=1.56, w=6.2)
    text_runs(s, 7.35, 1.72, 5.15, 1.6, [
        {"text": "U-кривая 2023 распрямилась: ", "size": 12.5,
         "color": DEEP},
        {"text": "single-needle — до 99% на полном окне 1 млн",
         "size": 12.5, "bold": True, "color": DEEP},
        {"text": " (Gemini Deep Think).", "size": 12.5, "color": DEEP},
    ], line_spacing=1.2)
    # ── НИЖНИЙ ЯРУС: работу со смыслом — нет
    ocean_box(s, 0.55, 3.58, 12.25, 2.85, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.85, 3.68, 11.7, 0.4, [
        {"text": "Работу со смыслом — нет ", "size": 14.5, "bold": True,
         "color": DEEP},
        {"text": "(NoLiMa, 2025: бенчмарк убрал буквальное совпадение слов)",
         "size": 12.5, "color": SLATE},
    ])
    add_image(s, ASSETS / "charts/s25-nolima.png", x=0.85, y=4.12, w=4.6)
    text_runs(s, 5.75, 4.25, 6.8, 1.7, [
        {"text": "Без лексических совпадений качество падает ", "size": 13,
         "color": DEEP},
        {"text": "ниже 50% от базового уже на 32K", "size": 13.5,
         "bold": True, "color": GOLD},
        {"text": " — у 11 из 13 моделей.", "size": 13, "bold": True,
         "color": DEEP},
        {"text": "32 тыс. токенов ≈ 3% заявленного окна флагмана; база — "
                 "точность той же модели на коротком контексте.",
         "size": 11.5, "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 8},
    ], line_spacing=1.2)
    gold_callout(s, 0.55, 6.55, 12.25, 0.72,
                 "1M окна ≠ 1M рассуждения. Окно — сколько модель может "
                 "прочитать; полезная длина — на скольких токенах она ещё "
                 "связывает факты.", size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s25"))


# ============================================================
# БАТЧ 2 — Раздел 4. Сэмплинг и генерация
# ============================================================
def build_s26a(p):
    section_divider(
        p, section_n=4, sub_title="Сэмплинг и генерация",
        frame_phrase="Как из распределения вероятностей рождается один "
                     "токен — и какими ручками этот выбор управляется",
        tag="4 разбора · 2 провала", active_stage={4, 5}, notes_id="s26a",
        passed={0, 1, 2, 3},
        frame_size=18,
        illus=ASSETS / "web/dice-wikimedia.jpg",
        illus_caption="Казино-кости Caesars Palace — Wikimedia Commons, "
                      "CC-BY-SA")


S26_BARS = [("яблоко", 0.32, True), ("пиццу", 0.19, False),
            ("салат", 0.14, False), ("булочку", 0.11, False),
            ("огурец", 0.08, False)]


def build_s26(p):
    """Распределение top-5 (нативные бары) + сэмплинг → один токен +
    footnote + gold callout."""
    s = blank(p)
    slide_title(s, "На каждом шаге — распределение вероятностей на все "
                   "токены словаря", size=24, h=0.9, y=0.35)
    # Левый box — бары
    ocean_box(s, 0.55, 1.5, 7.9, 3.7, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.9, 1.7, 7.2, 0.45, [
        {"text": "Контекст: ", "size": 15, "bold": True, "color": MID},
        {"text": "«Сегодня я съел …»", "size": 15.5, "bold": True,
         "color": DEEP, "font": FONT_MONO},
    ])
    text_box(s, 0.9, 2.2, 7.2, 0.35, "P(следующий токен):", size=13,
             color=SLATE, italic=True)
    yy = 2.62
    for lab, v, is_gold in S26_BARS:
        text_box(s, 0.9, yy, 1.35, 0.4, lab, size=13.5, bold=is_gold,
                 color=DEEP, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.RIGHT)
        bw = v / 0.35 * 4.6
        filled_rect(s, 2.4, yy + 0.05, bw, 0.32,
                    GOLD if is_gold else MID, radius=True, radius_adj=0.3)
        text_runs(s, 2.5 + bw, yy, 1.6, 0.4, [
            {"text": f"{v:.2f}".replace(".", ","), "size": 13.5,
             "bold": True, "color": GOLD if is_gold else DEEP},
            {"text": "  — максимум" if is_gold else "", "size": 11,
             "italic": True, "color": SLATE},
        ], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.5
    # Правый box — сэмплинг
    ocean_box(s, 8.75, 1.5, 4.05, 3.7, fill=WHITE, stroke=TEAL,
              stroke_pt=1.3)
    text_box(s, 9.0, 1.7, 3.55, 0.45, "Сэмплинг → один токен", size=15,
             bold=True, color=TEAL)
    text_box(s, 9.0, 2.25, 3.55, 0.7,
             "правило выбора одного токена из распределения — "
             "единственная ручка в ваших руках",
             size=12, color=DEEP, line_spacing=1.2)
    down_arrow(s, 10.6, 3.1, w=0.4, h=0.5, fill=TEAL)
    chip(s, 9.55, 3.75, 2.4, 0.55, "[ яблоко ]", fill=GOLD, color=DEEP,
         size=15, font=FONT_MONO)
    text_box(s, 9.0, 4.42, 3.55, 0.6, "выбран один — остальные кандидаты "
             "исчезли из ответа", size=11, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER, line_spacing=1.15)
    text_box(s, 0.55, 5.38, 12.25, 0.4,
             "Остальные ~200 000 токенов словаря — каждый < 0.05. Сумма "
             "всех вероятностей = 1. Числа иллюстративные.",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 5.85, 12.25, 1.0,
                 "Распределение — «настоящий» выход модели. Уверенный ответ "
                 "и галлюцинация до сэмплинга существовали одновременно — "
                 "как вероятностная масса; выбор сделала политика.",
                 size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s26"))


def _mini_bars(s, x, y, w, h, vals, *, gold_first=False):
    """Мини-распределение: вертикальные бары в области (x,y,w,h)."""
    n = len(vals)
    bw = w / (n * 1.5)
    gap = bw * 0.5
    vmax = max(vals) or 1.0
    xx = x + (w - (n * bw + (n - 1) * gap)) / 2
    for i, v in enumerate(vals):
        bh = max(v / vmax * h, 0.02)
        filled_rect(s, xx, y + h - bh, bw, bh,
                    GOLD if (gold_first and i == 0) else MID,
                    radius=False)
        xx += bw + gap
    plain_line(s, x, y + h + 0.02, x + w, y + h + 0.02, color=LIGHT,
               w_pt=1.2)


def build_s27(p):
    """v3.0: явная формула (argmax / P^(1/T)) + 3 панели распределений +
    callout «почему T влияет» + top-p/top-k + live-бейдж."""
    s = blank(p)
    slide_title(s, "Температура — делитель логитов: меняет остроту выбора, "
                   "не знания", size=25, h=0.55, y=0.32)
    # Формула — явно, в Ocean rounded box
    ocean_box(s, 0.55, 0.98, 12.25, 0.85, fill=SURFACE, stroke=MID,
              stroke_pt=1.4)
    text_runs(s, 0.85, 1.04, 11.7, 0.75, [
        {"text": "T = 0  ⇒  выбор = argmax P(токен)", "size": 14.5,
         "bold": True, "color": DEEP, "font": FONT_MONO},
        {"text": "      T > 0  ⇒  сэмплинг из P^(1/T)-образного "
                 "распределения", "size": 14.5, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": "логиты делятся на T до softmax: T<1 заостряет "
                 "распределение, T>1 выравнивает", "size": 11.5,
         "italic": True, "color": SLATE, "newpara": True,
         "space_before_pt": 4},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line_spacing=1.1)
    panels = [
        ("T → 0", "(argmax)", [1.0, 0.001, 0.001, 0.001, 0.001],
         "Выбор самого вероятного. Почти одинаковые ответы — «почти» "
         "разберём на следующем слайде.", False),
        ("T = 1", "(стандарт)", [0.32, 0.19, 0.14, 0.11, 0.08],
         "Сэмплирование пропорционально вероятностям модели.", True),
        ("T = 1.5", "(сглаживание)", [0.24, 0.19, 0.16, 0.14, 0.12],
         "Редкие токены получают реальные шансы — до бессвязности.",
         False),
    ]
    x = 0.55
    for tname, tsub, vals, caption, is_std in panels:
        if is_std:
            filled_rect(s, x, 2.0, 3.95, 2.95, WHITE, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.05)
        else:
            ocean_box(s, x, 2.0, 3.95, 2.95, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.3)
        text_runs(s, x + 0.25, 2.14, 3.45, 0.42, [
            {"text": tname, "size": 16, "bold": True, "color": DEEP,
             "font": FONT_MONO},
            {"text": "  " + tsub, "size": 13.5, "bold": True,
             "color": GOLD if is_std else MID},
        ])
        _mini_bars(s, x + 0.55, 2.62, 2.85, 1.15, vals, gold_first=True)
        text_box(s, x + 0.25, 4.02, 3.45, 0.85, caption, size=11,
                 color=DEEP, line_spacing=1.15)
        x += 4.15
    # Почему T влияет, хотя порядок токенов не меняется
    gold_callout(s, 0.55, 5.10, 12.25, 0.78,
                 "Почему T влияет, хотя порядок кандидатов не меняется: "
                 "выбор — не «взять топ-1», а случайная выборка "
                 "пропорционально вероятностям; T перекраивает сами "
                 "вероятности.", size=13, align=PP_ALIGN.CENTER)
    text_runs(s, 0.55, 6.00, 12.3, 0.42, [
        {"text": "top-p", "size": 13, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": " — отрез хвоста по вероятностной массе · ", "size": 12.5,
         "color": DEEP},
        {"text": "top-k", "size": 13, "bold": True, "color": MID,
         "font": FONT_MONO},
        {"text": " — по числу кандидатов · T под задачу: 0–0.3 код и "
                 "классификация, 0.7+ генерация.", "size": 12.5,
         "color": DEEP},
    ], align=PP_ALIGN.CENTER)
    ocean_box(s, 2.65, 6.52, 8.0, 0.55, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_runs(s, 2.85, 6.52, 7.6, 0.55, [
        {"text": "Живой прогон: ", "size": 12.5, "bold": True,
         "color": TEAL},
        {"text": "один запрос — 10 раз при T=0 и 10 раз при T=1.5.",
         "size": 12.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """T=0 ≠ детерминизм: экспонат 80/1000 + цепочка причины + 2 плашки +
    gold callout. Закрывает M4."""
    s = blank(p)
    slide_title(s, "T=0 не даёт детерминизма: 80 уникальных ответов из 1000",
                size=25, h=0.6)
    # Слева — экспонат (высота под контент, освобождает место мему снизу)
    ocean_box(s, 0.55, 1.35, 4.55, 2.15, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.8, 1.48, 4.05, 0.85, "80 / 1000", size=48, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, 0.85, 2.28, 3.95, 1.15,
             "уникальных вариантов ответа на идентичный запрос при T=0 — "
             "стандартный vLLM (открытый инференс-сервер; Thinking "
             "Machines Lab, сентябрь 2025)",
             size=11, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.12)
    # Реальный мем «Well yes, but actually no» — тот же шаблон, что s01,
    # другой текст: T=0 звучит как гарантия детерминизма, но нет. Помещаем
    # в зазор между экспонатом (низ 3.50) и рядом плашек (верх 4.75).
    meme_h = 1.08
    meme_w = meme_h * 1600 / 1218
    meme_x = 0.55 + (4.55 - meme_w) / 2
    meme_y = 3.60
    add_image(s, ASSETS / "web/well-yes-actually-no-template.jpg",
              x=meme_x, y=meme_y, h=meme_h)
    band_h = meme_h * 0.205
    text_box(s, meme_x + 0.08, meme_y + 0.02, meme_w - 0.16, band_h - 0.03,
             "T=0 = детерминизм?", size=10, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=0.95)
    # Справа — механизм
    text_runs(s, 5.4, 1.35, 7.4, 0.6, [
        {"text": "Причина — не «floating-point вообще», а отсутствие ",
         "size": 13.5, "bold": True, "color": DEEP},
        {"text": "batch-инвариантности ядер:", "size": 13.5, "bold": True,
         "color": MID},
    ], line_spacing=1.15)
    chain = [
        "Сервер группирует одновременные запросы разных пользователей в батчи",
        "Размер батча зависит от чужой нагрузки в эту миллисекунду",
        "Разный размер батча → разный порядок суммирования → младшие "
        "разряды другие",
        "Два близких кандидата в argmax → младший разряд решает выбор "
        "токена → авторегрессия разносит расхождение по всему ответу",
    ]
    yy = 2.0
    for i, step in enumerate(chain):
        h = 0.55 if i < 3 else 0.7
        ocean_box(s, 5.4, yy, 7.4, h, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, 5.6, yy, 7.0, h, [
            {"text": f"{i+1}. ", "size": 12, "bold": True, "color": TEAL},
            {"text": step, "size": 11.5, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        yy += h + 0.08
    # Две плашки
    filled_rect(s, 0.55, 4.75, 6.0, 1.0, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.3, radius=True, radius_adj=0.1)
    text_runs(s, 0.75, 4.82, 5.6, 0.9, [
        {"text": "Решение существует: ", "size": 12.5, "bold": True,
         "color": TEAL},
        {"text": "batch-инвариантные ядра — 1000/1000 побитово идентичны",
         "size": 12.5, "color": DEEP},
    ], line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)
    ocean_box(s, 6.8, 4.75, 6.0, 1.0, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.2)
    text_runs(s, 7.0, 4.82, 5.6, 0.9, [
        {"text": "Цена: ~35% пропускной способности", "size": 12.5,
         "bold": True, "color": DEEP},
        {"text": " → провайдеры не включают; seed у OpenAI — официально "
                 "«mostly deterministic», в основном детерминировано",
         "size": 12, "color": DEEP},
    ], line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 6.05, 12.25, 0.95,
                 "Получить гарантированно детерминированный ответ от "
                 "облачной LLM сегодня нельзя — стройте процессы с учётом "
                 "этого. Следствие: тесты — не на побитовом сравнении, а "
                 "семантически или по структуре.",
                 size=14, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """v3.0: таблица «параметр — диапазон — на что влияет — типичное
    значение» (5 ручек, effort — gold). Строка «заучивайте оси» убрана;
    budget_tokens/ошибка 400 — только в speaker notes."""
    s = blank(p)
    slide_title(s, "Ручки API: к случайности и длине добавилась глубина "
                   "рассуждения", size=24, h=0.6)
    # Таблица 6×4
    tx, ty, tw = 0.55, 1.45, 12.25
    col_ws = [2.9, 1.85, 4.1, 3.4]
    gtbl = s.shapes.add_table(6, 4, Inches(tx), Inches(ty), Inches(tw),
                              Inches(4.35))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, w in enumerate(col_ws):
        tbl.columns[ci].width = Inches(w)
    tbl.rows[0].height = Inches(0.45)
    for ri in range(1, 6):
        tbl.rows[ri].height = Inches(0.78)

    def cell(r, c, txt, *, size=12, bold=False, color=DEEP, fill=WHITE,
             mono=False, align=PP_ALIGN.LEFT):
        cl = tbl.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
        cl.margin_left = Inches(0.08); cl.margin_right = Inches(0.06)
        cl.margin_top = Inches(0.02); cl.margin_bottom = Inches(0.02)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cl.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        para.line_spacing = 1.08
        r_ = para.add_run(); r_.text = txt
        r_.font.name = FONT_MONO if mono else FONT_BODY
        r_.font.size = Pt(size); r_.font.bold = bold
        r_.font.color.rgb = color

    headers = ["Параметр", "Диапазон", "На что влияет",
               "Типичное значение"]
    for ci, htxt in enumerate(headers):
        cell(0, ci, htxt, size=12.5, bold=True, color=MID)
    rows_data = [
        ("temperature", "0–2", "Детерминизм ↔ хаос выбора токена",
         "0 классификация; 0.7–1.2 текст", False),
        ("top_p", "0.1–1", "Ширина хвоста кандидатов при сэмплинге",
         "0.9–0.95", False),
        ("max_tokens", "целое", "Жёсткий обрыв генерации — может прервать "
         "на середине", "под задачу, с запасом на JSON/код", False),
        ("effort", "none → xhigh", "Глубина внутреннего рассуждения — "
         "и его цена", "medium у большинства провайдеров", True),
        ("verbosity", "low → high", "Длина видимого ответа, независимо от "
         "глубины мышления", "medium", False),
    ]
    for ri, (name, rng, eff, typ, is_gold) in enumerate(rows_data):
        fill = GOLD_TINT if is_gold else (SURFACE if ri % 2 == 0 else WHITE)
        cell(ri + 1, 0, name, size=13, bold=True,
             color=DEEP if is_gold else MID, fill=fill, mono=True)
        cell(ri + 1, 1, rng, size=12, fill=fill, mono=True,
             align=PP_ALIGN.CENTER)
        cell(ri + 1, 2, eff, size=12, fill=fill, bold=is_gold)
        cell(ri + 1, 3, typ, size=11.5, fill=fill)
    text_box(s, 0.55, 5.95, 12.25, 0.35,
             "effort / reasoning_effort — новая ось 2026: OpenAI — шкала "
             "effort; Anthropic — адаптивное мышление; Gemini — thinking "
             "budget.", size=11.5, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.40, 12.25, 0.68,
                 "Что делать: начинайте настройку с temperature и effort — "
                 "две главные ручки; top_p / top_k / verbosity — тонкая "
                 "настройка поверх.", size=13.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """Constrained decoding: механика маскирования + ограничения +
    retrieval-пауза."""
    s = blank(p)
    slide_title(s, "Structured outputs: невалидные токены обнуляются прямо "
                   "в распределении", size=23, h=0.9, y=0.35)
    # Слева — механика
    ocean_box(s, 0.55, 1.5, 6.9, 4.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_runs(s, 0.85, 1.65, 6.3, 1.15, [
        {"text": "Constrained decoding: ", "size": 13, "bold": True,
         "color": MID},
        {"text": "JSON-схема компилируется в ", "size": 13, "color": DEEP},
        {"text": "конечный автомат над токенами", "size": 13, "bold": True,
         "color": DEEP},
        {"text": ". При потокенной генерации автомат отслеживает состояние "
                 "префикса и на каждом шаге МАСКИРУЕТ (обнуляет "
                 "вероятность) токены, ведущие к невалидному продолжению.",
         "size": 13, "color": DEEP},
    ], line_spacing=1.18)
    # мини-схема маскирования: 5 баров, 2 погашены
    vals = [(0.32, True), (0.19, False), (0.14, True), (0.11, False),
            (0.08, True)]
    xx = 1.45
    for v, valid in vals:
        bh = v / 0.35 * 1.0
        filled_rect(s, xx, 2.95 + (1.0 - bh), 0.6, bh,
                    MID if valid else SOFT_GREY, radius=False)
        if not valid:
            text_box(s, xx + 0.08, 2.82 + (1.0 - bh) - 0.28, 0.5, 0.35, "×",
                     size=17, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        xx += 1.0
    plain_line(s, 1.25, 3.98, 6.45, 3.98, color=LIGHT, w_pt=1.2)
    text_box(s, 1.25, 4.04, 5.4, 0.3,
             "автомат маскирует серые → вероятность 0",
             size=11, italic=True, color=SLATE)
    text_runs(s, 0.85, 4.48, 6.3, 1.1, [
        {"text": "Просьба «ответь строго в JSON» → ~80% валидных\n",
         "size": 13, "color": DEEP},
        {"text": "Strict mode → ", "size": 13.5, "bold": True,
         "color": DEEP, "newpara": True, "space_before_pt": 4},
        {"text": "100%", "size": 15, "bold": True, "color": GOLD},
        {"text": " — выход валиден по построению, а не «проверен "
                 "постфактум»", "size": 12.5, "color": DEEP},
    ], line_spacing=1.18)
    # Справа — ограничения
    ocean_box(s, 7.75, 1.5, 5.05, 4.25, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    text_box(s, 8.0, 1.65, 4.55, 0.4,
             "Ограничения — свойства компиляции", size=14, bold=True,
             color=MID)
    limits = [
        "Рекурсия через $ref — нет (дерево → плоский список с parent_id)",
        "Глубина вложенности — ≤ 5",
        "Первый запрос с новой схемой платит за компиляцию грамматики — "
        "до 10 с",
        "Гарантирован синтаксис, не смысл: валидировать значения "
        "по-прежнему вам",
    ]
    yy = 2.15
    for lim in limits:
        filled_rect(s, 8.0, yy, 4.55, 0.76, SURFACE, stroke=LIGHT,
                    stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(s, 8.15, yy + 0.04, 4.25, 0.68, lim, size=11, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        yy += 0.87
    ocean_box(s, 2.65, 6.05, 8.0, 0.65, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_runs(s, 2.85, 6.05, 7.6, 0.65, [
        {"text": "Вопрос залу: ", "size": 14, "bold": True, "color": TEAL},
        {"text": "почему именно 100%, а не 99.9?", "size": 14,
         "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """v3.0 рефрейм вокруг стоп-условий: цикл 5 шагов + возврат + карточки
    «как останавливается» / «как ломается» (вырожденные повторы, мем
    «горшочек, не вари!») + практика (repetition penalty, max_tokens)."""
    s = blank(p)
    slide_title(s, "Предсказали токен → дописали в контекст → предсказываем "
                   "следующий", size=21, h=0.55)
    steps = [
        ("1 · Текущий контекст",
         "промпт + история + запрос + уже сгенерированное", False),
        ("2 · Прямой проход",
         "токенизация → эмбеддинги → все слои внимания", True),
        ("3 · Распределение",
         "вероятности на все ~200 тыс. токенов", False),
        ("4 · Сэмплинг",
         "один токен — температура / top-p / схема", False),
        ("5 · Токен дописан",
         "в контекст — и цикл повторяется", False),
    ]
    bw, gap = 2.32, 0.24
    x0 = (SLIDE_W_IN - bw * 5 - gap * 4) / 2
    y0, bh = 1.72, 1.55
    for i, (head, desc, is_gold) in enumerate(steps):
        x = x0 + i * (bw + gap)
        if is_gold:
            filled_rect(s, x, y0, bw, bh, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.08)
        else:
            ocean_box(s, x, y0, bw, bh, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.3)
        text_box(s, x + 0.15, y0 + 0.10, bw - 0.3, 0.62, head, size=12.5,
                 bold=True, color=DEEP, line_spacing=1.08)
        text_box(s, x + 0.15, y0 + 0.72, bw - 0.3, 0.78, desc, size=10.5,
                 color=MID if not is_gold else DEEP, line_spacing=1.12)
        if i < 4:
            right_arrow(s, x + bw + 0.02, y0 + bh / 2 - 0.09, w=gap - 0.04,
                        h=0.18, fill=MID)
    chip(s, x0 - 0.15, y0 - 0.40, 0.9, 0.32, "вход", fill=GOLD, color=DEEP,
         size=11.5)
    # Возврат
    plain_line(s, x0 + 4 * (bw + gap) + bw / 2, y0 + bh,
               x0 + 4 * (bw + gap) + bw / 2, y0 + bh + 0.30, color=LIGHT,
               w_pt=2.0)
    left_arrow(s, x0 + bw / 2, y0 + bh + 0.24, w=4 * (bw + gap), h=0.18,
               fill=LIGHT)
    plain_line(s, x0 + bw / 2, y0 + bh, x0 + bw / 2, y0 + bh + 0.26,
               color=LIGHT, w_pt=2.0)
    text_box(s, x0 + 3.0, y0 + bh + 0.44, 6.0, 0.32,
             "⟲ возврат к шагу 1 — цикл повторяется", size=11.5,
             italic=True, color=MID, align=PP_ALIGN.CENTER)
    # Как останавливается / как ломается
    ocean_box(s, 0.55, 4.15, 5.95, 1.75, fill=WHITE, stroke=TEAL,
              stroke_pt=1.4)
    text_box(s, 0.85, 4.27, 5.4, 0.38, "Как цикл останавливается",
             size=14, bold=True, color=TEAL)
    text_runs(s, 0.85, 4.68, 5.4, 1.1, [
        {"text": "Специальный ", "size": 12.5, "color": DEEP},
        {"text": "стоп-токен", "size": 12.5, "bold": True, "color": DEEP},
        {"text": " — модель сама решает «ответ закончен» — или ",
         "size": 12.5, "color": DEEP},
        {"text": "max_tokens", "size": 12.5, "bold": True, "color": DEEP,
         "font": FONT_MONO},
        {"text": ": обрыв мгновенный, хоть на середине JSON-поля.",
         "size": 12.5, "color": DEEP},
    ], line_spacing=1.18)
    ocean_box(s, 6.8, 4.15, 6.0, 1.75, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 7.1, 4.27, 4.0, 0.38, "Как цикл ломается",
             size=14, bold=True, color=DEEP)
    text_box(s, 7.1, 4.68, 3.05, 1.1,
             "Вырожденный цикл повторов: модель «зацикливается» на одном "
             "токене или фразе и генерирует полотно повторов вместо "
             "остановки.", size=12, color=DEEP, line_spacing=1.15)
    # Реальная иллюстрация: «Горшочек каши» / «Sweet Porridge» (братья
    # Гримм) — Otto Ubbelohde, 1909, общественное достояние. Каша,
    # заполняющая деревню, — метафора цикла без стоп-условия.
    gsh_w = 1.55
    gsh_h = gsh_w * 324 / 640
    gsh_x = 10.55
    gsh_y = 4.30 + (1.1 - gsh_h) / 2
    add_image(s, ASSETS / "web/gorshochek-ubbelohde-1909.jpg",
              x=gsh_x, y=gsh_y, w=gsh_w)
    text_box(s, gsh_x - 0.12, gsh_y + gsh_h + 0.03, gsh_w + 0.24, 0.28,
             "«Горшочек каши» — Гримм, 1909", size=7.5, italic=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.05, 12.25, 0.68,
                 "Практика: repetition_penalty / frequency_penalty снижают "
                 "вероятность буквальных повторов; max_tokens — страховка "
                 "от бесконечного цикла.", size=13, align=PP_ALIGN.CENTER)
    text_box(s, 0.55, 6.82, 12.3, 0.4,
             "Каждый шаг — без состояния: вся «память» живёт в контексте, "
             "который подаётся целиком (KV-cache делает подачу дешёвой, не "
             "отменяя её логически).",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    """Reasoning-токены: раздвоенный выход + ценовые бары + границы +
    gold callout. Закрывает M5."""
    s = blank(p)
    slide_title(s, "Reasoning-токены не видны — но тарифицируются как output",
                size=25, h=0.6)
    # Верх — раздвоенный выход
    ocean_box(s, 0.55, 1.3, 12.25, 2.25, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    filled_rect(s, 0.95, 1.9, 2.7, 0.75, MID, radius=True, radius_adj=0.12)
    text_box(s, 1.0, 1.95, 2.6, 0.65, "авторегрессионный\nцикл  ⟲",
             size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    line_arrow(s, 3.7, 2.1, 4.4, 1.8, color=TEAL, w_pt=2.2)
    line_arrow(s, 3.7, 2.45, 4.4, 2.8, color=SLATE, w_pt=2.2)
    filled_rect(s, 4.5, 1.55, 4.6, 0.5, TEAL, radius=True, radius_adj=0.2)
    text_box(s, 4.65, 1.58, 4.3, 0.44, "видимый ответ", size=12.5,
             bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 4.5, 2.55, 7.7, 0.62, SOFT_GREY, radius=True,
                radius_adj=0.14)
    text_box(s, 4.7, 2.55, 6.2, 0.62,
             "черновик «для себя» — в ответ не попадает, в счёт попадает",
             size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.1)
    text_box(s, 11.05, 2.5, 1.15, 0.72, "×3–10", size=19, bold=True,
             color=GOLD, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text_box(s, 4.5, 3.24, 8.2, 0.3,
             "по ставке output-токенов, с учётом в max_tokens — без "
             "естественного потолка", size=11, italic=True, color=SLATE)
    # Низ слева — ценовые бары
    ocean_box(s, 0.55, 3.75, 5.95, 2.2, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.3)
    pb = [("o4-mini", 1, MID), ("o3", 5, MID), ("o3-pro", 18, GOLD)]
    yy = 3.95
    for lab, v, color in pb:
        text_box(s, 0.8, yy, 1.05, 0.32, lab, size=11.5, bold=True,
                 color=DEEP, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
        bw = v / 18 * 3.6
        filled_rect(s, 1.95, yy + 0.04, max(bw, 0.18), 0.24, color,
                    radius=True, radius_adj=0.35)
        text_box(s, 2.05 + max(bw, 0.18), yy, 0.9, 0.32, f"{v}×",
                 size=12, bold=True, color=GOLD if color == GOLD else DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.42
    text_runs(s, 0.8, 5.28, 5.5, 0.6, [
        {"text": "o3-pro: 3.6× дороже o3, 18× дороже o4-mini", "size": 11.5,
         "bold": True, "color": DEEP},
        {"text": " — при сопоставимой длине видимых ответов; разницу делает "
                 "объём рассуждения.", "size": 11, "color": SLATE},
    ], line_spacing=1.12)
    # Низ справа — 2 границы
    borders = [
        ("Управление:", " адаптивное мышление / effort вместо ручных "
         "бюджетов — удобно, но стоимость запроса стала менее "
         "предсказуемой"),
        ("«Ход рассуждений» в интерфейсе — пересказ", " (summarized), не "
         "сырые токены: строить аудит решений на «показанных мыслях» "
         "нельзя"),
    ]
    yy = 3.75
    for head, rest in borders:
        ocean_box(s, 6.75, yy, 6.05, 1.02, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, 6.95, yy + 0.06, 5.65, 0.9, [
            {"text": head, "size": 12, "bold": True, "color": MID},
            {"text": rest, "size": 11.5, "color": DEEP},
        ], line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.18
    gold_callout(s, 0.55, 6.15, 12.25, 0.85,
                 "Закладывайте в бюджет невидимую часть 2–5× видимого "
                 "ответа — и сверяйте по полю usage в ответе API, где "
                 "reasoning-токены видны как строка расхода.",
                 size=14, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s32"))


# ============================================================
# v3.0 — Раздел 5. Виды и размеры моделей (НОВЫЙ)
# ============================================================
def build_s33a(p):
    """NEW v3.0: divider нового Раздела 5. Вся лента конвейера в
    gold-подсветке (раздел логически «над» конвейером — классификация
    моделей, исполняющих его; НЕ отдельная стадия). Без gold-рамки
    (рамка — отличительный знак полного прохода на s35a)."""
    section_divider(
        p, section_n=5, sub_title="Виды и размеры моделей",
        frame_phrase="На чём запускаются модели разных размеров — и что "
                     "каждый класс умеет",
        tag="4 разбора", active_stage=set(), notes_id="s33a",
        bar_muted=True, frame_size=18,
        illus=ASSETS / "web/matryoshka-wikimedia.jpg",
        illus_caption="Русская матрёшка — Wikimedia Commons, CC-BY-SA "
                      "3.0 / GFDL")


S33_COLS = [
    ("s33-laptop", "Малые — до 8–10B", False, [
        ("Примеры: ", "Qwen3.8-4B / 8B, модели Llama-класса"),
        ("Железо: ", "ноутбук, смартфон, edge-устройство"),
        ("Мультимодальность: ", "обычно text-only или базовый vision"),
    ]),
    ("s33-gpu", "Средние — около 30B", False, [
        ("Примеры: ", "Muse Glimmer 30B — верх среднего класса"),
        ("Железо: ", "одна GPU 24–32 ГБ"),
        ("Мультимодальность: ", "часто есть vision"),
    ]),
    ("s33-server", "Большие — 70B+", False, [
        ("Примеры: ", "модели Llama-класса 70B, Qwen3.5-397B-A17B"),
        ("Железо: ", "мульти-GPU / сервер"),
        ("Мультимодальность: ", "как правило полная (текст + изображение, "
         "иногда аудио)"),
    ]),
    ("s33-cloud", "Гиганты-MoE — 400B+", True, [
        ("Примеры: ", "DeepSeek V4-Pro (1.6 трлн), Kimi K3 (2.8 трлн)"),
        ("Железо: ", "ТОЛЬКО облако или кластер — на потребительское "
         "железо не помещаются"),
        ("Мультимодальность: ", "полная, топ качество"),
    ]),
]


def build_s33(p):
    """NEW v3.0 (поз.37): классификация моделей по размеру — 4-колоночная
    матрица (иконка железа сверху каждого столбца), гиганты-MoE gold +
    callout про память как лимитирующий фактор."""
    s = blank(p)
    slide_title(s, "Четыре класса моделей по размеру — от ноутбука до "
                   "кластера", size=25, h=0.6)
    col_w, gap = 3.02, 0.14
    x0 = (SLIDE_W_IN - col_w * 4 - gap * 3) / 2
    y0, col_h = 1.28, 4.15
    for i, (icon, title, is_gold, rows) in enumerate(S33_COLS):
        x = x0 + i * (col_w + gap)
        if is_gold:
            filled_rect(s, x, y0, col_w, col_h, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.05)
        else:
            ocean_box(s, x, y0, col_w, col_h, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.4)
        add_image(s, ASSETS / f"icons/{icon}.png",
                  x=x + col_w / 2 - 0.27, y=y0 + 0.16, w=0.54)
        text_box(s, x + 0.1, y0 + 0.78, col_w - 0.2, 0.62, title,
                 size=14, bold=True, color=DEEP if is_gold else MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.05)
        runs = []
        for j, (head, rest) in enumerate(rows):
            runs.append({"text": head, "size": 12.5, "bold": True,
                         "color": TEAL if not is_gold else MID,
                         "newpara": j > 0, "space_before_pt": 14})
            runs.append({"text": rest, "size": 12.5, "color": DEEP,
                         "bold": is_gold and head.startswith("Железо")})
        text_runs(s, x + 0.2, y0 + 1.52, col_w - 0.4, col_h - 1.65, runs,
                  line_spacing=1.22)
    gold_callout(s, 0.55, 5.50, 12.25, 0.95,
                 "Лимитирующий фактор — объём памяти, не вычисления. Чем "
                 "крупнее модель, тем шире мультимодальность и качество — "
                 "и тем меньше шансов запустить её у себя.",
                 size=14, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s33"))


def build_s34(p):
    """Local vs cloud 2026: три категории; средняя — категориальная
    новость (gold)."""
    s = blank(p)
    slide_title(s, "«Открытые веса» перестали означать «локально "
                   "запускаемые»", size=26, h=0.6)
    cols = [
        ("Действительно локальные — до ~30B", None, [
            "Qwen3.8-27B (Apache 2.0, вход изображение+видео, окно "
            "262 тыс.), Muse Glimmer 30B",
            "Железо: RTX 5090 (32 ГБ) · Apple unified 64–128 ГБ",
            "Лимит — объём памяти, не вычисления"]),
        ("Открытые, но облачные гиганты", "открытая ≠ локальная", [
            "Kimi K3 — 2.8 трлн параметров, крупнейшая открытая модель",
            "DeepSeek V4-Pro — 1.6 трлн",
            "На потребительское железо не помещаются ни в каком виде"]),
        ("Закрытые API", None, [
            "Флагманское качество — по-прежнему здесь",
            "Плата за токен, данные через провайдера"]),
    ]
    x = 0.55
    for title, badge, bullets in cols:
        is_mid = badge is not None
        if is_mid:
            filled_rect(s, x, 1.5, 3.95, 4.0, WHITE, stroke=GOLD,
                        stroke_pt=2.2, radius=True, radius_adj=0.05)
            chip(s, x + 0.55, 1.28, 2.85, 0.4, badge, fill=GOLD, color=DEEP,
                 size=13)
        else:
            ocean_box(s, x, 1.5, 3.95, 4.0, fill=SURFACE, stroke=LIGHT,
                      stroke_pt=1.4)
        # Первый заголовок — одной строкой без переноса (v2.0.2 item 7)
        t_sz = 13.5 if "~30B" in title else 15
        text_box(s, x + 0.15, 1.82, 3.75, 0.75, title, size=t_sz, bold=True,
                 color=MID if not is_mid else DEEP, line_spacing=1.1)
        runs = []
        for i, b in enumerate(bullets):
            bold = b.startswith(("Лимит", "На потребительское"))
            runs.append({"text": "• " + b, "size": 12.5, "color": DEEP,
                         "bold": bold, "newpara": i > 0,
                         "space_before_pt": 8})
        text_runs(s, x + 0.25, 2.7, 3.45, 2.6, runs, line_spacing=1.2)
        x += 4.15
    text_box(s, 0.55, 5.85, 12.25, 0.55,
             "Причины локального выбора прежние: приватность данных · "
             "отсутствие платы за токен на объёмах · независимость от сети.",
             size=12.5, italic=True, color=MID, align=PP_ALIGN.CENTER,
             line_spacing=1.15)
    gold_callout(s, 1.4, 6.45, 10.55, 0.72,
                 "Что делать: задача умещается в ~30B и данные нельзя "
                 "выпускать за периметр — берите локальную модель; нужно "
                 "флагманское качество — берите закрытый API.",
                 size=13, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s34"))


# ============================================================
# v3.0 — Раздел 6. Финал (было «Раздел 5»)
# ============================================================
def build_s35a(p):
    section_divider(
        p, section_n=6, sub_title="Финал",
        frame_phrase="Конвейер собран целиком: сводка механизмов и границ — "
                     "и решение, когда LLM не тот инструмент",
        tag="6 разборов", active_stage=set(range(7)),
        notes_id="s35a", frame_bar=True, frame_size=18,
        illus=ASSETS / "web/this-is-fine-meme-fb.jpg",
        illus_caption="«This is fine» — K.C. Green, комикс Gunshow #648, "
                      "2013")


def build_s35(p):
    """v3.0 recap-фрейм: «Мы рассмотрели, как работает модель, — соберём
    картину»; плашка Chat-шаблонов → Glitch-токены (s07-каскад)."""
    s = blank(p)
    slide_title(s, "Мы рассмотрели, как работает модель, — соберём картину",
                size=24, h=0.6)
    text_box(s, 0.55, 1.02, 12.3, 0.4,
             "Новые темы не добавили стадий — они встроились в существующие.",
             size=14, italic=True, color=MID)
    stages = ["Токенизация", "Эмбеддинги", "Внимание", "Сэмплинг"]
    bw, gap, y0, bh = 2.28, 0.3, 3.3, 0.85
    x0 = 0.62
    centers = []
    for i, st in enumerate(stages):
        x = x0 + i * (bw + gap)
        centers.append(x + bw / 2)
        ocean_box(s, x, y0, bw, bh, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.5)
        text_box(s, x, y0 + 0.08, bw, bh - 0.16, st, size=14.5, bold=True,
                 color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            right_arrow(s, x + bw + 0.03, y0 + bh / 2 - 0.09, w=gap - 0.06,
                        h=0.18, fill=MID)
    rx = x0 + 4 * (bw + gap)
    filled_rect(s, rx, y0 + 0.08, 1.85, bh - 0.16, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.8, radius=True, radius_adj=0.15)
    text_runs(s, rx, y0 + 0.12, 1.85, bh - 0.24, [
        {"text": "следующий", "size": 12, "bold": True, "color": DEEP,
         "align": PP_ALIGN.CENTER},
        {"text": "токен", "size": 12, "bold": True, "color": DEEP,
         "newpara": True, "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    right_arrow(s, rx - gap + 0.03, y0 + bh / 2 - 0.09, w=gap - 0.06,
                h=0.18, fill=MID)
    # loop назад
    plain_line(s, rx + 0.92, y0 + bh - 0.08, rx + 0.92, y0 + bh + 0.32,
               color=LIGHT, w_pt=1.8)
    left_arrow(s, centers[0], y0 + bh + 0.26, w=rx + 0.92 - centers[0],
               h=0.16, fill=LIGHT)
    text_box(s, 4.2, y0 + bh + 0.46, 5.0, 0.3, "⟲ цикл — токен дописан в "
             "контекст", size=11, italic=True, color=MID,
             align=PP_ALIGN.CENTER)
    # Плашки-наложения: 2 сверху, 2 снизу
    overlays_top = [
        (0, 0.62, 3.35, "Glitch-токены", " — свойство словаря на стадии "
         "токенизации"),
        (2, 4.9, 4.5, "KV-cache", " — внутри внимания; prompt caching — "
         "надстройка над ним на границе запросов"),
    ]
    for stage_i, px, pw, head, rest in overlays_top:
        filled_rect(s, px, 1.55, pw, 0.95, TEAL_TINT, stroke=GOLD,
                    stroke_pt=1.5, radius=True, radius_adj=0.1)
        text_runs(s, px + 0.18, 1.60, pw - 0.36, 0.85, [
            {"text": head, "size": 12.5, "bold": True, "color": DEEP},
            {"text": rest, "size": 11.5, "color": DEEP},
        ], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
        plain_line(s, px + pw / 2, 2.5, centers[stage_i], y0, color=TEAL,
                   w_pt=1.4, dash=4)
    overlays_bottom = [
        (3, 7.6, 4.0, "Structured outputs", " — фильтр на стадии сэмплинга",
         centers[3]),
        (None, 2.9, 4.3, "Reasoning-токены", " — тот же цикл; часть выхода "
         "помечена «черновик»", centers[0] + 1.0),
    ]
    for stage_i, px, pw, head, rest, tx_ in overlays_bottom:
        filled_rect(s, px, 5.0, pw, 0.95, TEAL_TINT, stroke=GOLD,
                    stroke_pt=1.5, radius=True, radius_adj=0.1)
        text_runs(s, px + 0.18, 5.05, pw - 0.36, 0.85, [
            {"text": head, "size": 12.5, "bold": True, "color": DEEP},
            {"text": rest, "size": 11.5, "color": DEEP},
        ], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
        plain_line(s, px + pw / 2, 5.0,
                   tx_ if stage_i is None else centers[stage_i],
                   y0 + bh + (0.3 if stage_i is None else 0),
                   color=TEAL, w_pt=1.4, dash=4)
    gold_callout(s, 0.55, 6.25, 12.25, 0.85,
                 "Конвейер — диагностическое дерево: по симптому почти "
                 "всегда угадывается стадия-виновник. Спрашивайте: «на какой "
                 "стадии это происходит?»", size=14.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s35"))


def build_s36(p):
    """Ландшафт 2026: 2 лагеря + IMO-экспонат + ценовая шкала."""
    s = blank(p)
    slide_title(s, "Сентябрь 2026: качество сблизилось — цены разошлись на "
                   "три порядка", size=24, h=0.9, y=0.35)
    # Колонки — реальные логотипы компаний (LobeHub icons-static-svg,
    # recolor в Ocean palette) как визуальный якорь перед каждой строкой.
    frontier = [
        ("openai", "OpenAI: GPT-5.6 — Luna → Terra → Sol"),
        ("anthropic", "Anthropic: Claude Fable 5 · Opus 5"),
        ("google", "Google: Gemini 3.5 Pro (окно 2 млн, Deep Think)"),
        ("xai", "xAI: Grok 4.3"),
    ]
    open_w = [
        ("deepseek", "DeepSeek V4 (Pro 1.6 трлн / Flash 284 млрд)"),
        ("qwen", "Qwen 3.8-Max — первая открытая из линейки Max"),
        (None, "Kimi K2.6 (1 трлн) · Kimi K3 (2.8 трлн — крупнейшая "
               "открытая)"),
    ]
    for x, title, items in [(0.55, "Передний край (закрытые веса)",
                             frontier),
                            (6.8, "Открытые веса", open_w)]:
        ocean_box(s, x, 1.5, 6.0, 2.4, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.4)
        text_box(s, x + 0.25, 1.65, 5.5, 0.4, title, size=15, bold=True,
                 color=MID)
        row_y = 2.14
        row_h = 1.7 / len(items)
        for logo, it in items:
            bold = "крупнейшая открытая" in it
            if logo:
                add_image(s, ASSETS / f"icons/logos-web/{logo}.png",
                          x=x + 0.25, y=row_y + row_h / 2 - 0.14, h=0.28)
                tx0 = x + 0.62
            else:
                tx0 = x + 0.25
            text_box(s, tx0, row_y, x + 5.75 - tx0, row_h, it, size=12.5,
                     bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.1)
            row_y += row_h
    # IMO экспонат
    filled_rect(s, 0.55, 4.05, 12.25, 0.92, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.8, radius=True, radius_adj=0.12)
    text_box(s, 0.85, 4.13, 11.7, 0.4,
             "IMO 2026: шесть моделей — абсолютные 42/42. Из 666 "
             "участников-людей — семеро.", size=14.5, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, 0.85, 4.55, 11.7, 0.38,
             "Те же системы ошибаются в подсчёте букв слова cranberry — "
             "«рваный интеллект» как рабочая характеристика.",
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)
    # Ценовая шкала
    filled_rect(s, 1.0, 5.55, 11.3, 0.2, SOFT_GREY, radius=True,
                radius_adj=0.5)
    filled_rect(s, 1.0, 5.55, 2.3, 0.2, LIGHT, radius=True, radius_adj=0.5)
    filled_rect(s, 10.0, 5.55, 2.3, 0.2, DEEP, radius=True, radius_adj=0.5)
    text_box(s, 1.0, 5.20, 4.5, 0.32, "пол рынка $0.03–0.2 / млн токенов",
             size=12, bold=True, color=MID)
    text_box(s, 8.3, 5.20, 4.0, 0.32, "премиум $10 вход / $50 выход",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.RIGHT)
    text_runs(s, 0.55, 5.98, 12.3, 0.4, [
        {"text": "Kimi K2.6 ≈ GPT-5.5 на защищённом SWE-bench Pro — ",
         "size": 13, "bold": True, "color": DEEP},
        {"text": "при цене на ~80% ниже", "size": 13, "bold": True,
         "color": GOLD},
    ], align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 6.55, 12.25, 0.68,
                 "Что делать: пересматривайте выбор модели регулярно — "
                 "ландшафт живёт месяцами, а не годами.",
                 size=13, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s36"))


def build_s37(p):
    """Бенчмарки: 3 сюжета-карточки + gold callout. Закрывает M6."""
    s = blank(p)
    slide_title(s, "Бенчмарки: контаминация, подгонка — и модели, которые "
                   "жульничают сами", size=24, h=0.9, y=0.35)
    # v2.0.2 (item 10): текст сюжетов ≤2 строк (детали в notes);
    # цифры 87.6 vs 57 — крупным контрастным стат-блоком.
    cards = [
        ("1", "Контаминация: заучено, а не умеет", 8.2, [
            {"text": "SWE-bench: публичные репозитории (Verified) vs "
                     "приватные базы (Pro) — ", "size": 12.5, "color": DEEP},
            {"text": "разрыв и есть величина заучивания.", "size": 12.5,
             "bold": True, "color": DEEP},
            {"text": " OpenAI в 2026 перестал публиковать Verified.",
             "size": 12.5, "color": DEEP}]),
        ("2", "Подгонка под метрику", 11.0, [
            {"text": "Llama 4 Maverick: на Chatbot Arena — специальная "
                     "версия, Elo 1417; публичная модель — ", "size": 12.5,
             "color": DEEP},
            {"text": "места 32–35", "size": 12.5, "bold": True,
             "color": DEEP},
            {"text": ". Ян Лекун: результаты «слегка подтасованы».",
             "size": 12.5, "color": DEEP}]),
        ("3", "Модели жульничают сами", 11.0, [
            {"text": "UK AI Security Institute: ", "size": 12.5,
             "color": DEEP},
            {"text": "все 5", "size": 12.5, "bold": True, "color": DEEP},
            {"text": " передовых моделей пытались обмануть процедуру "
                     "оценки; одна модель OpenAI ", "size": 12.5,
             "color": DEEP},
            {"text": "вышла из песочницы и взломала производственные "
                     "серверы Hugging Face", "size": 12.5,
             "bold": True, "color": DEEP},
            {"text": ".", "size": 12.5, "color": DEEP}]),
    ]
    yy = 1.52
    for num, head, body_w, runs in cards:
        ocean_box(s, 0.55, yy, 12.25, 1.42, fill=SURFACE if num != "3"
                  else WHITE, stroke=LIGHT, stroke_pt=1.3)
        text_box(s, 0.8, yy + 0.1, 0.55, 0.6, num, size=26, bold=True,
                 color=LIGHT)
        text_box(s, 1.45, yy + 0.12, 10.9, 0.38, head, size=14.5,
                 bold=True, color=MID)
        text_runs(s, 1.45, yy + 0.52, body_w, 0.85, runs, line_spacing=1.16)
        if num == "1":
            # Стат-блок: 87.6% (gold) vs 57% — крупно, контрастно
            text_runs(s, 9.85, yy + 0.16, 2.8, 0.75, [
                {"text": "87.6%", "size": 24, "bold": True, "color": GOLD},
                {"text": " vs ", "size": 14, "color": SLATE},
                {"text": "57%", "size": 24, "bold": True, "color": DEEP},
            ], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            text_box(s, 8.85, yy + 0.94, 3.8, 0.4,
                     "заявка производителя vs защищённый · средний ~25%",
                     size=10.5, italic=True, color=SLATE,
                     align=PP_ALIGN.RIGHT)
        yy += 1.55
    gold_callout(s, 0.55, 6.25, 12.25, 0.8,
                 "Бенчмарки сужают список кандидатов. Выбирает — "
                 "собственный оценочный набор: 30–50 примеров из ваших "
                 "реальных задач.", size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s37"))


# v3.0: «Подведём итоги» — сводная таблица «механизм → граница → что
# делать» по всем реально рассмотренным темам (без M-кодов, без чек-листа).
S38_ROWS = [
    ("Токенизация",
     "Модель видит токены, не буквы; вирусный фикс не переносится на класс",
     "Тестируйте «cranberry своей предметки»; счёт — в инструмент"),
    ("Внимание и роль",
     "Роль — токены с весом; смещает стиль и фокус, не качество",
     "Роль — инструмент тона/фокуса; для фактов — давайте данные"),
    ("KV-cache и кэш промптов",
     "Экономит только на повторном префиксе; токен в начале рушит всё после",
     "Стабильное — в начало, переменное — в конец; следите за cache hit"),
    ("Контекстное окно",
     "Заявленное ≠ полезное: 11 из 13 моделей теряют половину уже на 32K",
     "Выбирайте по бенчмаркам без лексических подсказок"),
    ("Детерминизм при T=0",
     "Ядра не batch-инвариантны: чужая нагрузка меняет ваш ответ",
     "Тесты — не на побитовом сравнении; стройте процессы с учётом"),
    ("Reasoning-токены",
     "Тарифицируются как output, без естественного потолка",
     "Закладывайте бюджет отдельно; effort/verbosity — явно"),
    ("Структурированный вывод",
     "Валидность по построению — но не содержательное качество",
     "Не перегружайте схему; значения валидируете вы"),
    ("Бенчмарки",
     "Контаминация, подгонка витрин, жульничество моделей",
     "Свой оценочный набор; лидерборды — ориентир, не гарантия"),
    ("Размеры моделей",
     "«Открытые веса» ≠ «локально запускаемые»: гиганты — только облако",
     "Класс модели — по задаче и железу, не по лицензии весов"),
]


def build_s38(p):
    """v3.0 «Подведём итоги»: сводная таблица «механизм → граница → что
    делать» (9 строк, без M-кодов, без переклички с s01)."""
    s = blank(p)
    slide_title(s, "Подведём итоги: по каждому механизму — граница и что "
                   "с ней делать", size=24, h=0.6)
    tx, ty, tw = 0.55, 1.18, 12.25
    n_rows = len(S38_ROWS) + 1
    gtbl = s.shapes.add_table(n_rows, 3, Inches(tx), Inches(ty),
                              Inches(tw), Inches(4.95))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, w in enumerate([2.55, 5.1, 4.6]):
        tbl.columns[ci].width = Inches(w)
    tbl.rows[0].height = Inches(0.36)
    for ri in range(1, n_rows):
        tbl.rows[ri].height = Inches(0.51)

    def cell(r, c, txt, *, size=10.5, bold=False, color=DEEP, fill=WHITE,
             align=PP_ALIGN.LEFT):
        cl = tbl.cell(r, c)
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
        cl.margin_left = Inches(0.07); cl.margin_right = Inches(0.05)
        cl.margin_top = Inches(0.01); cl.margin_bottom = Inches(0.01)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cl.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        para.line_spacing = 1.02
        r_ = para.add_run(); r_.text = txt
        r_.font.name = FONT_BODY; r_.font.size = Pt(size)
        r_.font.bold = bold
        r_.font.color.rgb = color

    cell(0, 0, "Механизм", size=12, bold=True, color=MID)
    cell(0, 1, "Граница", size=12, bold=True, color=MID)
    cell(0, 2, "Что делать", size=12, bold=True, color=MID)
    for ri, (mech, bound, act) in enumerate(S38_ROWS):
        fill = SURFACE if ri % 2 == 0 else WHITE
        cell(ri + 1, 0, mech, size=11, bold=True, color=MID, fill=fill)
        cell(ri + 1, 1, bound, size=10.5, fill=fill)
        cell(ri + 1, 2, act, size=10.5, color=DEEP, fill=fill)
    gold_callout(s, 0.55, 6.32, 12.25, 0.75,
                 "Знать инструмент — значит знать его границы. Каждый "
                 "механизм работает — но не безгранично.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s38"))


def build_s39(p):
    """Decision tree «когда не LLM» + лестница эскалации."""
    s = blank(p)
    slide_title(s, "Когда не LLM — и когда не топ-LLM", size=26, h=0.6)
    text_box(s, 0.55, 1.15, 7.5, 0.4, "Когда LLM — не тот инструмент:",
             size=15, bold=True, color=MID)
    branches = [
        ("Классификация на фиксированных категориях с тысячами размеченных "
         "примеров", "классический ML: дешевле, быстрее, воспроизводимо — "
         "а LLM ещё и недетерминирована при T=0"),
        ("Объяснимость перед регулятором", "прозрачные классические методы"),
        ("Отклик < 100 мс (антифрод, устройства без сети)",
         "специализированная малая модель"),
        ("Точные посимвольные и арифметические операции", "код, не модель"),
    ]
    yy = 1.62
    for cond, ans in branches:
        ocean_box(s, 0.55, yy, 7.6, 0.92, fill=SURFACE, stroke=LIGHT,
                  stroke_pt=1.2)
        text_runs(s, 0.75, yy + 0.05, 7.2, 0.84, [
            {"text": cond, "size": 12, "bold": True, "color": DEEP},
            {"text": "  →  ", "size": 12.5, "bold": True, "color": MID},
            {"text": ans, "size": 12, "color": TEAL, "bold": True},
        ], line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.02
    filled_rect(s, 0.55, yy, 7.6, 0.85, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.4, radius=True, radius_adj=0.1)
    text_runs(s, 0.75, yy + 0.04, 7.2, 0.77, [
        {"text": "Иначе", "size": 12.5, "bold": True, "color": TEAL},
        {"text": " — обработка языка, гибкие форматы, многошаговое "
                 "рассуждение, генерация → ", "size": 12, "color": DEEP},
        {"text": "LLM применима и часто оптимальна", "size": 12,
         "bold": True, "color": DEEP},
    ], line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
    # Справа — лестница
    ocean_box(s, 8.4, 1.15, 4.4, 5.55, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 8.65, 1.3, 3.9, 0.4, "…и не всегда топ-LLM", size=15,
             bold=True, color=DEEP)
    filled_rect(s, 10.3, 2.25, 2.25, 0.8, DEEP, radius=True,
                radius_adj=0.12)
    text_box(s, 10.42, 2.3, 2.0, 0.7, "10% сложных →\nпремиум $10/млн",
             size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    filled_rect(s, 8.65, 3.2, 3.9, 0.9, MID, radius=True, radius_adj=0.1)
    text_box(s, 8.8, 3.25, 3.6, 0.8, "90% запросов →\nмодель за $0.20/млн",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.08)
    text_runs(s, 8.65, 4.55, 3.9, 1.7, [
        {"text": "Миллиард токенов/мес:", "size": 12.5, "bold": True,
         "color": DEEP},
        {"text": "$10 000", "size": 14, "bold": True, "color": DEEP,
         "newpara": True, "space_before_pt": 8},
        {"text": " целиком на премиуме", "size": 12, "color": DEEP},
        {"text": "vs  ", "size": 12.5, "color": SLATE, "newpara": True,
         "space_before_pt": 6},
        {"text": "$1 180", "size": 16, "bold": True, "color": GOLD},
        {"text": " с маршрутизацией", "size": 12, "color": DEEP},
    ], line_spacing=1.15)
    speaker_notes(s, load_notes("s39"))


def build_s40(p):
    """Корреляция ≠ причинность: человек vs модель + gold callout."""
    s = blank(p)
    slide_title(s, "Внимание усваивает корреляцию, не причинность", size=26,
                h=0.6)
    # Человек
    ocean_box(s, 0.55, 1.45, 5.95, 3.55, fill=SURFACE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 0.85, 1.62, 5.35, 0.4, "Человек", size=16.5, bold=True,
             color=MID)
    text_runs(s, 0.85, 2.1, 5.35, 0.85, [
        {"text": "«X произошло, потому что Y»", "size": 14, "bold": True,
         "color": DEEP},
        {"text": " — модель механизмов мира", "size": 13.5, "color": DEEP},
    ], line_spacing=1.2)
    hlevels = [("ассоциация", "✓"), ("вмешательство", "✓"),
               ("контрфактуальность", "✓")]
    yy = 3.0
    for lab, mark in hlevels:
        filled_rect(s, 0.85, yy, 5.35, 0.52, TEAL_TINT, stroke=TEAL,
                    stroke_pt=1.2, radius=True, radius_adj=0.2)
        text_runs(s, 1.05, yy, 5.0, 0.52, [
            {"text": mark + "  ", "size": 13, "bold": True, "color": TEAL},
            {"text": lab, "size": 13, "bold": True, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.62
    # Модель
    ocean_box(s, 6.85, 1.45, 5.95, 3.55, fill=WHITE, stroke=LIGHT,
              stroke_pt=1.4)
    text_box(s, 7.15, 1.62, 5.35, 0.4, "Модель (через внимание)",
             size=16.5, bold=True, color=DEEP)
    text_runs(s, 7.15, 2.1, 5.35, 0.85, [
        {"text": "«X следует за Y в текстах»", "size": 14, "bold": True,
         "color": DEEP},
        {"text": " — «потому что» для модели — частотный паттерн, не "
                 "указание на механизм мира", "size": 12.5, "color": DEEP},
    ], line_spacing=1.18)
    mlevels = [
        ("ассоциация — сильна", "✓", TEAL_TINT, TEAL, TEAL),
        ("вмешательство — только похожие на описанные в корпусе", "~",
         WHITE, LIGHT, MID),
        ("контрфактуальность — систематически нет", "×", SOFT_GREY, SLATE,
         SLATE),
    ]
    yy = 3.0
    for lab, mark, fill, stroke, mcol in mlevels:
        filled_rect(s, 7.15, yy, 5.35, 0.52, fill, stroke=stroke,
                    stroke_pt=1.2, radius=True, radius_adj=0.2)
        text_runs(s, 7.35, yy, 5.0, 0.52, [
            {"text": mark + "  ", "size": 13, "bold": True, "color": mcol},
            {"text": lab, "size": 11.5, "bold": True, "color": DEEP},
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        yy += 0.62
    gold_callout(s, 0.55, 5.35, 12.25, 0.95,
                 "Там, где от модели ждут каузальных выводов, человек в "
                 "контуре — архитектурное требование, а не вежливая "
                 "оговорка.", size=15.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s40"))


def build_s41(p):
    """Мост к Лекции 3 (v3.1, #183 round 3 — fix композиции): мост теперь
    ГОРИЗОНТАЛЬНОЙ ПОЛОСОЙ под заголовком, ПОЛНОСТЬЮ читаемый (арка +
    тросы + опоры), НЕ перекрыт карточками — раньше карточки лежали
    поверх моста, виден был только обрубок вертикальной полосы. Картинка
    предварительно обрезана по контенту (убрано пустое верхнее поле PNG:
    784/1024 исходной высоты), поэтому силуэт моста заполняет всю полосу
    без потерь. 4 карточки-концепта — компактный 2×2 grid ниже полосы."""
    s = blank(p)
    text_runs(s, 0.55, 0.30, 12.3, 0.55, [
        {"text": "Лекция 3: ", "size": 24, "bold": True, "color": GOLD},
        {"text": "как модель выходит за пределы контекста", "size": 24,
         "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # Hero-полоса: мост, обрезанный по контенту (1984×784) — читается
    # целиком: арка, тросы, обе опоры, конечные узлы слева/справа.
    bridge_w = 7.4
    bridge_h = bridge_w * 784 / 1984
    bridge_x = (SLIDE_W_IN - bridge_w) / 2
    bridge_y = 0.90
    add_image(s, ASSETS / "illustrations/s41-bridge-lec3-crop.png",
              x=bridge_x, y=bridge_y, w=bridge_w)
    cards = [
        ("s41-search", "RAG",
         "Семантический поиск по вашей базе → найденные фрагменты в "
         "контекст.",
         "Якорь: сходство ≠ релевантность.", TEAL),
        ("s41-settings", "Инструменты / вызов функций",
         "Модель генерирует структурированный вызов → внешняя система "
         "исполняет.",
         "Якорь: надёжность формата даёт structured outputs.", TEAL),
        ("s41-plug", "MCP",
         "Открытый протокол подключения инструментов.",
         "Якорь: стабильный префикс → кэш промптов, экономика агента.",
         GOLD),
        ("s41-refresh-cw", "Агентный цикл",
         "Действие → наблюдение → коррекция.",
         "Якорь: агент читает внешний контент → prompt injection.", TEAL),
    ]
    card_w, card_h = 5.75, 1.62
    gap_x, gap_y = 0.75, 0.14
    x0 = (SLIDE_W_IN - card_w * 2 - gap_x) / 2
    y0 = bridge_y + bridge_h + 0.14
    for i, (icon, title, body, anchor, acol) in enumerate(cards):
        col, row = i % 2, i // 2
        x = x0 + col * (card_w + gap_x)
        y = y0 + row * (card_h + gap_y)
        ocean_box(s, x, y, card_w, card_h, fill=WHITE, stroke=LIGHT,
                  stroke_pt=1.4)
        add_image(s, ASSETS / f"icons/{icon}.png", x=x + 0.18, y=y + 0.14,
                  w=0.34)
        text_box(s, x + 0.62, y + 0.11, card_w - 0.78, 0.36, title,
                 size=13.5, bold=True, color=DEEP)
        text_box(s, x + 0.22, y + 0.50, card_w - 0.44, 0.55, body, size=10.5,
                 color=DEEP, line_spacing=1.08)
        text_box(s, x + 0.22, y + 1.06, card_w - 0.44, 0.50, anchor,
                 size=9.5, italic=True,
                 color=acol if acol == GOLD else TEAL, line_spacing=1.08)
    speaker_notes(s, load_notes("s41"))


def build_s42(p):
    """Q&A minimal (паттерн Lec-1 s31): surface-фон, Q&A 140pt, Спасибо."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, 0.55, 1.7, 12.3, 2.6, "Q&A", size=140, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.55, 4.5, 12.3, 0.8, "Спасибо", size=36, italic=True,
             color=MID, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s42"))


# ============================================================
# Build — полный дек v3.0: 47 слайдов, 7 разделов
# ============================================================
BUILDERS = [
    # Раздел 0. Введение
    ("s01", build_s01), ("s02", build_s02), ("s02a", build_s02a),
    ("s03", build_s03), ("s04", build_s04), ("s04b", build_s04b),
    # Раздел 1. Токенизация (s07 удалён — тема → Лекция 3)
    ("s05a", build_s05a), ("s05", build_s05), ("s06", build_s06),
    ("s08", build_s08), ("s09", build_s09),
    ("s10", build_s10), ("s11", build_s11),
    # Раздел 2. Эмбеддинги
    ("s12a", build_s12a), ("s12", build_s12), ("s13", build_s13),
    ("s14", build_s14), ("s15", build_s15), ("s17", build_s17),
    # Раздел 3. Внимание (v3.0: KV-cache сразу после Q/K/V, роль после кэша)
    ("s18a", build_s18a), ("s18", build_s18), ("s19", build_s19),
    ("s21", build_s21), ("s22", build_s22), ("s20", build_s20),
    ("s23", build_s23), ("s25", build_s25),
    # Раздел 4. Сэмплинг и генерация
    ("s26a", build_s26a), ("s26", build_s26), ("s27", build_s27),
    ("s28", build_s28), ("s29", build_s29), ("s30", build_s30),
    ("s31", build_s31), ("s32", build_s32),
    # Раздел 5. Виды и размеры моделей (НОВЫЙ)
    ("s33a", build_s33a), ("s33", build_s33), ("s34", build_s34),
    ("s36", build_s36), ("s37", build_s37),
    # Раздел 6. Финал
    ("s35a", build_s35a), ("s35", build_s35), ("s38", build_s38),
    ("s39", build_s39), ("s40", build_s40), ("s41", build_s41),
    ("s42", build_s42),
]


def main():
    p = setup_pres()
    print(f"Building {len(BUILDERS)} slides (full deck v3.0)…")
    for sid, fn in BUILDERS:
        try:
            fn(p)
            print(f"  {sid} OK")
        except Exception as e:
            print(f"  {sid} FAIL: {type(e).__name__}: {e}")
            raise
    # Знаменатель — общее число слайдов в деке (47, включая s42 Q&A),
    # даже хотя сам s42 идёт без footer/номера (паттерн qa_minimal).
    num_total = len(BUILDERS)
    for i, ((sid, _fn), slide) in enumerate(zip(BUILDERS, p.slides)):
        if sid == "s42":
            continue  # Q&A — паттерн qa_minimal: без footer/номера страницы
        page_number(slide, f"{i + 1}/{num_total}")
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"\nSaved: {OUT}  ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
