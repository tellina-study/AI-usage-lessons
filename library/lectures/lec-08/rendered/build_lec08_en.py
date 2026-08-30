"""
Build script for Лекции 8 «AI в креативных индустриях и медиа» (39 slides).

Phase 5 deck v1 — Media-heavy ≥80% (≥33/39 slides with embedded media).

Source-of-truth: deck.yaml + chapter.md + slides/*.md (39 files with readable
speaker notes 150-300 words).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal
(#028090) secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke
#1C7293 1.5pt).

Canvas: 13.333" × 7.5" (16:9). Pacing ≈ 75 мин.

Build via: python3 build_lec08.py — generates lec-08.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
GREEN_OK  = RGBColor(0x2E, 0x8B, 0x57)
RED_WARN  = RGBColor(0xC0, 0x39, 0x2B)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path("/home/harness/harness-projects/256/.worktrees/folder-288/publish-8a63bf98/library/lectures/lec-08")
ASSETS = ROOT / "assets"
SLIDES_DIR = ROOT / "slides-en"
OUT = ROOT / "rendered/lec-08-en.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    if not Path(path).exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path)
            img_w, img_h = img.size
            img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        img_ratio = img_w / img_h
        box_ratio = w / h
        if img_ratio > box_ratio:
            actual_h = w / img_ratio
            y_offset = (h - actual_h) / 2
            slide.shapes.add_picture(str(path), Inches(x), Inches(y + y_offset),
                                     width=Inches(w))
        else:
            actual_w = h * img_ratio
            x_offset = (w - actual_w) / 2
            slide.shapes.add_picture(str(path), Inches(x + x_offset), Inches(y),
                                     height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True):
    box = filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                      radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.08, w=w - 0.4, h=h - 0.16, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


def lesson_box(slide, x, y, w, h, text):
    """«УРОК ДЛЯ ИНЖЕНЕРА» gold-tint Ocean rounded box."""
    box = ocean_box(slide, x, y, w, h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.8)
    text_box(slide, x=x + 0.18, y=y + 0.10, w=w - 0.36, h=0.32,
             text="LESSON FOR THE ENGINEER",
             size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT, line_spacing=1.0)
    text_box(slide, x=x + 0.18, y=y + 0.45, w=w - 0.36, h=h - 0.55,
             text=text,
             size=14, bold=True, color=DEEP, align=PP_ALIGN.LEFT, line_spacing=1.30)


def speaker_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def load_notes(slide_id):
    """Extract Speaker notes block from slide markdown (дословно copy)."""
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    notes_match = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                             md, re.DOTALL)
    notes = notes_match.group(1).strip() if notes_match else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Section divider — unified template (6 cards)
# ============================================================
NAV_SECTIONS = [
    ("0", "Opening\n+ keystone axis",   "3-tenses axis\n+ 3 families"),
    ("1", "AI ADDED",                   "new\ncapabilities"),
    ("2", "AI CHANGED",                 "cost · speed\n· professions"),
    ("3", "AI BROKE",                   "12 failure\ncases"),
    ("4", "AI NOT NEEDED",              "4 criteria\nfor refusal"),
    ("5", "What to do",                 "5-question\nchecklist"),
]


def build_section_divider(p, here_idx, title, frame_phrase, notes_slide_id):
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Big section number outline left
    text_box(s, x=0.55, y=1.0, w=4.0, h=4.5,
             text=str(here_idx),
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    # Section label
    text_box(s, x=4.7, y=1.6, w=8.0, h=0.5,
             text="SECTION",
             size=18, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Section title
    text_box(s, x=4.7, y=2.15, w=8.0, h=1.8,
             text=title,
             size=44, bold=True, color=DEEP,
             align=PP_ALIGN.LEFT, line_spacing=1.15)
    # Frame phrase
    filled_rect(s, 4.7, 3.95, 0.04, 0.55, fill=TEAL)
    text_box(s, x=4.85, y=3.95, w=7.9, h=1.0,
             text=frame_phrase,
             size=16, italic=False, color=MID,
             align=PP_ALIGN.LEFT, line_spacing=1.30)
    # Progress bar — 6 cards
    bar_y = 5.85
    n_cells = 6
    total_w = 12.3
    gap = 0.12
    cell_w = (total_w - gap * (n_cells - 1)) / n_cells
    start_x = 0.55
    card_h = 0.95
    for i, (num, sec_title, desc) in enumerate(NAV_SECTIONS):
        x = start_x + i * (cell_w + gap)
        is_here = (i == here_idx)
        if is_here:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=2.0)
            num_color = WHITE
            title_color = WHITE
        elif i < here_idx:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=TEAL_TINT, stroke=TEAL, stroke_pt=1.2)
            num_color = TEAL
            title_color = MID
        else:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=WHITE, stroke=LIGHT, stroke_pt=1.0)
            num_color = LIGHT
            title_color = SLATE
        text_box(s, x=x, y=bar_y + 0.08, w=cell_w, h=0.40, text=num,
                 size=22, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=x + 0.04, y=bar_y + 0.46, w=cell_w - 0.08, h=card_h - 0.50,
                 text=sec_title, size=9, bold=is_here, color=title_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.15)
    speaker_notes(s, load_notes(notes_slide_id))


# ============================================================
# Slide builders — 39 slides
# ============================================================

def build_s01(p):
    """Ice-breaker с hero-image (Sora 2 woolly mammoths) — instant emotional engagement
    + foreshadow keystone «AI добавил → изменил → сломал»."""
    s = blank(p)
    # ---- HERO IMAGE LEFT (Sora 2 mammoths — iconic AI-video frame) ----
    # Image area: 6.6×3.71 (16:9 aspect preserved) ≈ 24.5 sq in. Ocean motif обрамление.
    ocean_box(s, 0.45, 0.55, 6.70, 3.95, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, str(ASSETS / "screenshots/s01-sora-mammoths.png"),
              x=0.55, y=0.65, w=6.50, h=3.75)
    # Attribution chip
    text_box(s, x=0.55, y=4.55, w=6.5, h=0.30,
             text="OpenAI Sora · text-to-video frame from the prompt \"woolly mammoths\" · 2024",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # ---- ASSERTION + COST-COLLAPSE LEFT-BOTTOM ----
    text_box(s, x=0.55, y=4.95, w=6.6, h=1.05,
             text="AI produces a production-quality artifact in seconds, with no special skills.",
             size=20, bold=True, color=DEEP, line_spacing=1.15)
    # Cost-collapse compact 2-row
    text_runs(s, 0.55, 6.10, 6.6, 0.40, [
        {"text": "Music: ", "size": 13, "color": DEEP, "bold": True},
        {"text": "composer + a week + $500-2000  →  ", "size": 13, "color": SLATE, "italic": True},
        {"text": "Suno · 30 sec · $0", "size": 13, "color": GOLD, "bold": True},
    ], line_spacing=1.25)
    text_runs(s, 0.55, 6.50, 6.6, 0.40, [
        {"text": "Photo: ", "size": 13, "color": DEEP, "bold": True},
        {"text": "freelance $50-200 + 1-3 days  →  ", "size": 13, "color": SLATE, "italic": True},
        {"text": "Firefly · 5 sec · $0", "size": 13, "color": GOLD, "bold": True},
    ], line_spacing=1.25)
    # ---- RIGHT: DEMO CARD ----
    ocean_box(s, 7.55, 0.55, 5.30, 5.95)
    text_box(s, x=7.75, y=0.75, w=4.90, h=0.40,
             text="IN-BROWSER DEMO",
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, x=7.75, y=1.20, w=4.90, h=0.50,
             text="Let's generate one right now",
             size=20, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    # Primary card — Suno
    filled_rect(s, 7.85, 1.95, 4.70, 1.55, MID, radius=True, radius_adj=0.10)
    text_box(s, x=8.00, y=2.05, w=4.40, h=0.32,
             text="PRIMARY · AUDIO",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=2.42, w=4.40, h=0.55,
             text="suno.com/create",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=3.00, w=4.40, h=0.50,
             text="topic + genre + language → 30-sec track",
             size=11, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Fallback card — Firefly
    filled_rect(s, 7.85, 3.65, 4.70, 1.55, LIGHT, radius=True, radius_adj=0.10)
    text_box(s, x=8.00, y=3.75, w=4.40, h=0.32,
             text="BACKUP · IMAGE",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=4.12, w=4.40, h=0.55,
             text="firefly.adobe.com",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=4.70, w=4.40, h=0.50,
             text="text prompt → photo in 5 sec",
             size=11, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
    # QR placeholder
    filled_rect(s, 11.65, 5.40, 0.85, 0.85, WHITE, stroke=DEEP, stroke_pt=1.5)
    text_box(s, x=11.65, y=5.40, w=0.85, h=0.85, text="QR",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=7.85, y=5.50, w=3.60, h=0.70,
             text="Open the URL on your phone →",
             size=10, italic=True, color=DEEP, line_spacing=1.20,
             anchor=MSO_ANCHOR.MIDDLE)
    # Footer — Suno / Firefly URLs
    text_box(s, x=0.55, y=7.05, w=12.3, h=0.35,
             text="suno.com · firefly.adobe.com · free trial access · 30 sec",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Cover slide."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Big outline "08"
    text_box(s, x=8.0, y=2.0, w=5.3, h=5.5, text="08",
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    # Lecture label
    text_box(s, x=0.7, y=1.0, w=6.5, h=0.55, text="LECTURE",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.72, 1.55, 0.7, 0.05, fill=TEAL)
    # Main title
    text_box(s, x=0.7, y=2.0, w=8.5, h=2.6,
             text="AI in the creative\nindustries and media",
             size=54, bold=True, color=DEEP, line_spacing=1.05,
             align=PP_ALIGN.LEFT)
    # Subtitle
    filled_rect(s, 0.7, 5.45, 0.05, 0.6, fill=TEAL)
    text_box(s, x=0.95, y=5.45, w=10.5, h=0.7,
             text="What AI added, changed, and broke — and where to say \"no\".",
             size=20, color=MID, italic=False, align=PP_ALIGN.LEFT,
             line_spacing=1.25)
    # Gold highlight chip — student-facing tagline
    chip(s, 0.7, 6.5, 3.6, 0.45, "Lecture 08 · 75 min · 39 slides",
         fill=GOLD_TINT, stroke=GOLD, color=DEEP, size=11, bold=True)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    """Central question."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=0.5, w=12.3, h=0.7,
             text="CENTRAL QUESTION OF THE LECTURE",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
             line_spacing=1.0)
    # Big Ocean box with the question
    ocean_box(s, 1.2, 1.5, 10.9, 4.5)
    text_runs(s, 1.5, 1.8, 10.3, 4.0, [
        {"text": "What has AI done ", "size": 36, "color": GOLD, "bold": True},
        {"text": "to the creative industry by 2026 —", "size": 36, "color": DEEP, "bold": True},
        {"newpara": True, "text": "and where should an engineer reasonably say", "size": 36, "color": DEEP, "bold": True},
        {"newpara": True, "text": "\"AI is not needed here\"?", "size": 36, "color": TEAL, "bold": True},
    ], line_spacing=1.30, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Anchor below box
    text_box(s, x=0.55, y=6.3, w=12.3, h=0.5,
             text="Two-part by design — both parts matter equally.",
             size=15, italic=True, color=MID, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """Lecture map — 6-card horizontal дорожная карта."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.7,
             text="Lecture map — 6 sections + Q&A",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="6 sections · 75 min · 1 checklist to take away",
             size=16, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 6-card horizontal дорожная карта
    bar_y = 2.2
    n_cells = 6
    total_w = 12.3
    gap = 0.18
    cell_w = (total_w - gap * (n_cells - 1)) / n_cells
    start_x = 0.55
    card_h = 4.2
    here_idx = 0  # Currently in Section 0
    for i, (num, sec_title, desc) in enumerate(NAV_SECTIONS):
        x = start_x + i * (cell_w + gap)
        is_here = (i == here_idx)
        if is_here:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=2.0)
            num_color = WHITE
            title_color = WHITE
            desc_color = WHITE
        else:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=WHITE, stroke=LIGHT, stroke_pt=1.2)
            num_color = LIGHT
            title_color = DEEP
            desc_color = SLATE
        text_box(s, x=x, y=bar_y + 0.3, w=cell_w, h=1.5, text=num,
                 size=72, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=x + 0.10, y=bar_y + 2.0, w=cell_w - 0.20, h=1.0,
                 text=sec_title, size=14, bold=True, color=title_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.20)
        text_box(s, x=x + 0.10, y=bar_y + 3.0, w=cell_w - 0.20, h=1.0,
                 text=desc, size=10, color=desc_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)
    # Q&A separately
    chip(s, 5.5, 6.7, 2.5, 0.5, "+ Q&A",
         fill=TEAL_TINT, stroke=TEAL, color=DEEP, size=14, bold=True)
    speaker_notes(s, load_notes("s04"))


def build_s05(p):
    """KEYSTONE — AI добавил → изменил → сломал."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Main title
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.7,
             text="AI added → changed → broke",
             size=40, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.25, w=12.3, h=0.5,
             text="Three tenses of one process — each generation of creative tools passes through them in months.",
             size=16, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 timing bands — equal width, stacked vertically
    band_x = 0.55
    band_w = 12.3
    band_y_start = 2.1
    band_h = 1.45
    band_gap = 0.15
    bands = [
        ("ADDED", "New capabilities — things that technologically did not exist before.",
         "Sora 2 · 25 sec text-to-video · sync audio. Midjourney Omni Reference. ElevenLabs — voice clone from 1 min.",
         MID, WHITE),
        ("CHANGED", "Cost · speed · professions — the new economics of the creative industries.",
         "Cost collapse of 100×–10,000× across asset types. Concept art: days → seconds. Upwork +70% year over year in AI/ML.",
         LIGHT, WHITE),
        ("BROKE", "A new class of failures and legal debt.",
         "RIAA v. Suno, Sony's millions. Arup CFO deepfake $25.6M. SI — fake authors. NYT v. OpenAI — 20M logs.",
         GOLD, DEEP),
    ]
    for i, (label, subtitle, example, fill_color, txt_color) in enumerate(bands):
        y = band_y_start + i * (band_h + band_gap)
        # Color band
        filled_rect(s, band_x, y, 2.5, band_h, fill_color, radius=True, radius_adj=0.10)
        text_box(s, x=band_x + 0.20, y=y + 0.10, w=2.2, h=0.5, text=str(i+1),
                 size=22, bold=True, color=txt_color, align=PP_ALIGN.LEFT,
                 line_spacing=1.0)
        text_box(s, x=band_x + 0.20, y=y + 0.60, w=2.2, h=0.7,
                 text=label, size=22, bold=True, color=txt_color,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        # Subtitle + example
        ocean_box(s, band_x + 2.7, y, band_w - 2.7, band_h)
        text_box(s, x=band_x + 2.85, y=y + 0.15, w=band_w - 2.95, h=0.5,
                 text=subtitle, size=15, bold=True, color=DEEP,
                 align=PP_ALIGN.LEFT, line_spacing=1.20)
        text_box(s, x=band_x + 2.85, y=y + 0.65, w=band_w - 2.95, h=0.75,
                 text=example, size=12, italic=True, color=SLATE,
                 align=PP_ALIGN.LEFT, line_spacing=1.30)
    # Anchor footer
    text_box(s, x=0.55, y=7.0, w=12.3, h=0.35,
             text="This is an axis of three tenses of one process, not three parallel categories.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s05"))


def build_s05a(p):
    """3 families: diffusion / latent video transformer / neural audio."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.7,
             text="3 families of generative media models",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Mental model: each family → its own fundamental limitations, independent of implementation quality.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 cards equal width
    card_y = 2.0
    card_h = 4.7
    total_w = 12.3
    gap = 0.20
    card_w = (total_w - gap * 2) / 3
    start_x = 0.55
    families = [
        ("1. Diffusion",
         "noise → reversal → image",
         "Stable Diffusion · Midjourney · Flux · DALL-E · Imagen · Adobe Firefly",
         "\"Commercial safety\" depends on the training corpus, not on the architecture.",
         MID),
        ("2. Latent video transformers",
         "latent space + temporal consistency",
         "Sora 2 · Veo 3.1 · Runway · Kling 3.0",
         "Sora's 25-sec limit — cost grows linearly + temporal drift after ~25 sec.",
         LIGHT),
        ("3. Neural audio synthesis",
         "autoregressive + diffusion",
         "Suno · Udio · ElevenLabs · Stable Audio",
         "Voice clone from 1 min — fine-tuning a pretrained foundation model, not from scratch.",
         GOLD),
    ]
    for i, (title, principle, tools, consequence, accent) in enumerate(families):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        # Header band
        filled_rect(s, x, card_y, card_w, 0.55, accent,
                    radius=True, radius_adj=0.10)
        title_color = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=card_y + 0.10, w=card_w - 0.30, h=0.40,
                 text=title, size=17, bold=True, color=title_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Principle
        text_box(s, x=x + 0.18, y=card_y + 0.75, w=card_w - 0.36, h=0.4,
                 text="PRINCIPLE", size=10, bold=True, color=TEAL,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 1.05, w=card_w - 0.36, h=0.5,
                 text=principle, size=14, bold=True, color=DEEP, italic=True,
                 align=PP_ALIGN.LEFT, line_spacing=1.25)
        # Tools
        text_box(s, x=x + 0.18, y=card_y + 1.75, w=card_w - 0.36, h=0.4,
                 text="TOOLS 2026", size=10, bold=True, color=TEAL,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 2.05, w=card_w - 0.36, h=1.1,
                 text=tools, size=12, color=SLATE,
                 align=PP_ALIGN.LEFT, line_spacing=1.40)
        # Consequence
        text_box(s, x=x + 0.18, y=card_y + 3.25, w=card_w - 0.36, h=0.4,
                 text="ENGINEERING CONSEQUENCE", size=10, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 3.55, w=card_w - 0.36, h=1.05,
                 text=consequence, size=12, color=DEEP, bold=True,
                 align=PP_ALIGN.LEFT, line_spacing=1.30)
    speaker_notes(s, load_notes("s05a"))


# Section dividers handled by build_section_divider helper
def build_s06(p):
    build_section_divider(p, here_idx=1, title="AI ADDED",
                          frame_phrase="New capabilities: text-to-video · character preservation across generations · voice cloning · world models.",
                          notes_slide_id="s06")


def build_s07(p):
    """Text-to-video 2026 — 3-card comparison."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Text-to-video 2026 — production quality",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="3 flagship models define the state of the industry",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3-card comparison
    card_y = 2.0
    card_h = 4.0
    total_w = 8.0
    gap = 0.18
    card_w = (total_w - gap * 2) / 3
    start_x = 0.55
    models = [
        ("Sora 2", "OpenAI", "25 sec · 1080p · sync audio",
         "$0.10/sec 720p · cameos", MID, "openai.com/index/sora-2/"),
        ("Veo 3.1", "Google", "8 sec · 720p/1080p · native audio",
         "$0.05/sec Lite", LIGHT, "Google AI Ultra"),
        ("Kling 3.0", "Kuaishou", "15 sec · 4K · 60 fps",
         "ELO #1 (1243)", GOLD, "Feb 5, 2026 · 60M creators"),
    ]
    for i, (name, vendor, specs, price, accent, url) in enumerate(models):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        # Accent header
        filled_rect(s, x, card_y, card_w, 0.5, accent, radius=True, radius_adj=0.10)
        name_color = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.12, y=card_y + 0.06, w=card_w - 0.24, h=0.4,
                 text=name, size=18, bold=True, color=name_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.18, y=card_y + 0.65, w=card_w - 0.36, h=0.4,
                 text=vendor, size=11, bold=True, color=TEAL, italic=True)
        text_box(s, x=x + 0.18, y=card_y + 1.10, w=card_w - 0.36, h=1.0,
                 text=specs, size=13, color=DEEP, bold=True, line_spacing=1.35)
        text_box(s, x=x + 0.18, y=card_y + 2.20, w=card_w - 0.36, h=1.0,
                 text=price, size=13, color=GOLD, bold=True, line_spacing=1.30)
        text_box(s, x=x + 0.18, y=card_y + 3.20, w=card_w - 0.36, h=0.7,
                 text=url, size=10, italic=True, color=SLATE, line_spacing=1.20)
    # Right side: real Sora кадр из демо-релиза (woolly mammoths walking through snow,
    # OpenAI Sora YouTube channel · эталонная demo from launch)
    ocean_box(s, 8.75, 2.0, 4.20, 4.0)
    text_box(s, x=8.95, y=2.10, w=3.8, h=0.35,
             text="OPENAI SORA · LAUNCH DEMO · YOUTUBE",
             size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_image(s, str(ASSETS / "screenshots/s07-sora2-frame.png"),
              x=8.92, y=2.50, w=3.86, h=3.20)
    text_box(s, x=8.95, y=5.75, w=3.8, h=0.30,
             text="openai.com/index/sora-2/  ·  $0.10/sec",
             size=10, italic=True, color=LIGHT, line_spacing=1.15)
    # Мини-провал block
    lesson_box(s, 0.55, 6.2, 12.3, 1.0,
               "Mini-failure: 25 sec is not a film. A cinematic pipeline is assembled from short blocks under human direction (Lionsgate × Runway = augmentation, not replacement).")
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """Сохранение персонажа."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Character preservation: cameos and Omni Reference",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Multi-scene narrative became possible. Without it, AI video does not work for storytelling.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 cards vertical (left) — compressed to fit alongside lesson box
    cards_x = 0.55
    cards_w = 6.0
    card_h = 1.18
    card_gap = 0.10
    cards = [
        ("Sora 2 cameos", "OpenAI + Disney partnership worth $1B+",
         "A character is registered — called by name in prompts across multiple scenes."),
        ("Midjourney Omni Reference (v7)",
         "Image-to-image accuracy 60% → 85%+",
         "A character reference image preserves face/clothing/pose proportions."),
        ("Runway Gen-4 Director Mode",
         "Scripting across scenes via structured objects",
         "Character + locations + motion patterns as constants across all scenes."),
    ]
    for i, (title, subtitle, body) in enumerate(cards):
        y = 2.0 + i * (card_h + card_gap)
        ocean_box(s, cards_x, y, cards_w, card_h)
        text_box(s, x=cards_x + 0.18, y=y + 0.08, w=cards_w - 0.36, h=0.32,
                 text=title, size=14, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=cards_x + 0.18, y=y + 0.40, w=cards_w - 0.36, h=0.28,
                 text=subtitle, size=10, italic=True, color=TEAL, line_spacing=1.0)
        text_box(s, x=cards_x + 0.18, y=y + 0.68, w=cards_w - 0.36, h=0.48,
                 text=body, size=11, color=SLATE, line_spacing=1.20)
    # Right: real Midjourney референс персонажа grid («рыцарь в лесу» + «старик» с шляпой) —
    # 8 generated images showing character preservation across cw=0 vs cw=100 (consistency weight)
    ocean_box(s, 7.0, 2.0, 5.95, 3.95)
    text_box(s, x=7.15, y=2.10, w=5.6, h=0.35,
             text="MIDJOURNEY · CHARACTER REFERENCE · \"knight in a forest\" × \"old man\"",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s08-character-grid.png"),
              x=7.10, y=2.45, w=5.78, h=3.40)
    # Lesson box — positioned just below 3-card stack & grid
    lesson_box(s, 0.55, 6.05, 12.3, 1.20,
               "Drift between scenes appears after 5-10 scenes. Continuity supervisor — a new profession in the creative pipeline.")
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """Клонирование голоса + дубляж на несколько языков."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Voice cloning + dubbing into multiple languages",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="ElevenLabs: voice clone from 1 min of audio → 32+ languages. Production use in the corporate sector.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: real ElevenLabs официальная обложка (from elevenlabs.io/cover.png)
    ocean_box(s, 0.55, 2.0, 6.0, 4.0)
    text_box(s, x=0.75, y=2.10, w=5.6, h=0.35,
             text="ELEVENLABS · OFFICIAL · elevenlabs.io",
             size=12, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_image(s, str(ASSETS / "screenshots/s09-elevenlabs.png"),
              x=0.75, y=2.45, w=5.60, h=2.10)
    # Список голосов — concise, soundalike risk shown as Teal teal
    voices = [
        ("Voice 1", "Multilingual · expressive", MID),
        ("Voice 2", "Russian · neutral", MID),
        ("Voice 3", "EN-UK · narrator", LIGHT),
        ("Voice 4", "ScarJo-like soundalike", TEAL),
    ]
    for i, (vname, desc, badge_c) in enumerate(voices):
        y = 4.65 + i * 0.32
        filled_rect(s, 0.75, y, 5.60, 0.28, SURFACE, stroke=LIGHT, stroke_pt=0.5,
                    radius=True, radius_adj=0.20)
        filled_rect(s, 0.82, y + 0.05, 0.18, 0.18, badge_c, radius=True, radius_adj=0.5)
        text_box(s, x=1.10, y=y, w=5.20, h=0.28,
                 text=f"{vname}  ·  {desc}",
                 size=10, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # Right top: обвал стоимости arrow
    ocean_box(s, 6.85, 2.0, 6.10, 1.8)
    text_box(s, x=7.05, y=2.15, w=5.7, h=0.4,
             text="COST COLLAPSE: dubbing per language",
             size=12, bold=True, color=TEAL)
    text_runs(s, 7.05, 2.6, 5.7, 0.8, [
        {"text": "$50-500", "size": 24, "color": SLATE, "bold": True, "italic": True},
        {"text": "  →  ", "size": 26, "color": DEEP, "bold": True},
        {"text": "<$1", "size": 32, "color": GOLD, "bold": True},
        {"text": "  /min", "size": 18, "color": DEEP, "bold": True},
    ], line_spacing=1.10)
    text_box(s, x=7.05, y=3.40, w=5.7, h=0.35,
             text="Long-form: weeks → minutes (Dubbing Studio, 29 languages)",
             size=11, italic=True, color=LIGHT)
    # Right middle: productionrs
    text_box(s, x=6.85, y=3.95, w=6.10, h=0.4,
             text="PRODUCTION USE",
             size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    chip(s, 6.85, 4.35, 1.95, 0.45, "Deutsche Telekom",
         fill=MID, color=WHITE, size=12)
    chip(s, 8.95, 4.35, 1.85, 0.45, "Klarna",
         fill=MID, color=WHITE, size=12)
    chip(s, 10.95, 4.35, 2.00, 0.45, "Multi-language",
         fill=LIGHT, color=WHITE, size=11)
    # Мини-провал ScarJo
    lesson_box(s, 6.85, 5.05, 6.10, 1.95,
               "MINI-FAILURE: ScarJo v. OpenAI \"Sky\" (May 2024). OpenAI pulled the voice within a week — formally without a lawsuit. A de-facto win for the right of publicity: a voice clone requires explicit consent, even if it is technologically \"merely similar\".")
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    """Genie 3 world models."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="World models — Genie 3 (DeepMind)",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Text → playable 3D world @ 24 fps. This is NOT video generation — it is a simulated environment.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: real Genie 3 official 9-frame gameplay grid from DeepMind blog
    # (volcano, jellyfish, eagle, Japan, waterfall, Venice, wingsuit, alley + WASD controls)
    ocean_box(s, 0.55, 2.0, 7.0, 4.5)
    text_box(s, x=0.75, y=2.10, w=6.6, h=0.35,
             text="GENIE 3 · DEEPMIND · OFFICIAL DEMO",
             size=12, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s10-genie3-world.png"),
              x=0.70, y=2.50, w=6.70, h=3.55)
    text_box(s, x=0.75, y=6.10, w=6.6, h=0.30,
             text="deepmind.google/blog/genie-3-a-new-frontier-for-world-models",
             size=11, italic=True, color=LIGHT)
    # Right: 3 metric chips
    text_box(s, x=7.85, y=2.0, w=5.1, h=0.4,
             text="CHARACTERISTICS",
             size=12, bold=True, color=TEAL)
    metrics = [
        ("text → playable 3D world", MID),
        ("24 fps · real time", LIGHT),
        ("720p · several minutes", GOLD),
    ]
    for i, (mtext, color) in enumerate(metrics):
        y = 2.5 + i * 0.85
        filled_rect(s, 7.85, y, 5.1, 0.7, color, radius=True, radius_adj=0.20)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=7.85, y=y, w=5.1, h=0.7,
                 text=mtext, size=18, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    # Anti-hype — shifted up + shortened to fit
    lesson_box(s, 7.85, 5.20, 5.1, 1.80,
               "This is not a video generator — it is a simulated environment. Production use: so far edge cases (game prototypes, location scouting). The frontier is still ahead.")
    speaker_notes(s, load_notes("s10"))


def build_s10a(p):
    """Russian context — локальное удобство vs frontier."""
    s = blank(p)
    text_box(s, x=0.55, y=0.30, w=12.3, h=0.65,
             text="Russian context: local convenience vs the frontier",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.0, w=12.3, h=0.4,
             text="Russian GenAI is functional for the end user + mass-market premium; but not at the frontier in video and music. Structurally (capex/data), not ideologically.",
             size=12, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: real Russian AI image-gen ecosystem screenshot (Шедеврум web + mobile interfaces
    # с реальными AI-сгенерированными работами: осьминог-с-арбузом, крокодил-робот, дом-в-облаках)
    ocean_box(s, 0.55, 1.7, 5.5, 3.15)
    text_box(s, x=0.65, y=1.78, w=5.3, h=0.35,
             text="SHEDEVRUM + KANDINSKY · OFFICIAL INTERFACES",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s10a-kandinsky-vs-kling.png"),
              x=0.62, y=2.15, w=5.40, h=2.45)
    text_box(s, x=0.65, y=4.62, w=5.3, h=0.30,
             text="Source: appleinsider.ru comparison of Shedevrum vs Kandinsky · 2023",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right: 4-card RU landscape — Ocean palette only (no red/green)
    grid_x = 6.40
    grid_y = 1.7
    grid_w = 6.55
    grid_h = 3.0
    cell_w = (grid_w - 0.15) / 2
    cell_h = (grid_h - 0.15) / 2
    areas = [
        ("Images", "Kandinsky 6.0 (MoE, free via GigaChat)\nYandex Shedevrum · YandexART 2.7",
         "competitive", TEAL),
        ("Video", "Kandinsky 5.0 (Apache 2.0)\nLags behind Sora/Veo/Kling",
         "structural gap", GOLD),
        ("Audio", "SymFormer · SaluteSpeech\nYandex SpeechKit",
         "below ElevenLabs", MID),
        ("Legal landscape", "Mintsifry (Russian Ministry of Digital Development) 03/18/2026\nTDM* · content labeling · 09/01/2027",
         "in progress", LIGHT),
    ]
    for i, (atitle, abody, atag, atag_color) in enumerate(areas):
        row = i // 2; col = i % 2
        x = grid_x + col * (cell_w + 0.15)
        y = grid_y + row * (cell_h + 0.15)
        ocean_box(s, x, y, cell_w, cell_h)
        text_box(s, x=x + 0.12, y=y + 0.08, w=cell_w - 0.24, h=0.30,
                 text=atitle, size=12, bold=True, color=DEEP)
        text_box(s, x=x + 0.12, y=y + 0.40, w=cell_w - 0.24, h=0.65,
                 text=abody, size=10, color=SLATE, line_spacing=1.30)
        chip_text_c = DEEP if atag_color == GOLD else WHITE
        chip(s, x + 0.12, y + cell_h - 0.45, cell_w - 0.24, 0.35, atag,
             fill=atag_color, color=chip_text_c, size=10, bold=True)
    # Inline glossary footer (TDM)
    text_box(s, x=6.40, y=4.78, w=6.55, h=0.30,
             text="* TDM = Text & Data Mining (research exception in law)",
             size=9, italic=True, color=LIGHT, line_spacing=1.0)
    # Lesson box
    lesson_box(s, 0.55, 5.10, 12.3, 1.85,
               "Local convenience (free · RU prompts · no VPN · rubles · legal perimeter) ≠ frontier-level in video and music. The concentration of R&D in the US/China is structural (GPU capex, video datasets), not ideological.")
    speaker_notes(s, load_notes("s10a"))


def build_s11(p):
    """Personalisation at scale + Adobe Firefly."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Personalization at scale + production use in Hollywood",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Adobe Firefly: $400M direct revenue 2024-25. Lionsgate × Runway — augmentation, not replacement.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: Adobe Firefly коллаж логотипов + Lionsgate quote
    ocean_box(s, 0.55, 2.0, 7.0, 4.5)
    text_box(s, x=0.70, y=2.15, w=6.7, h=0.4,
             text="ADOBE FIREFLY · ENTERPRISE CLIENTS",
             size=11, bold=True, color=TEAL)
    # Logos as text chips
    логотипы = ["Deloitte", "Tapestry", "Paramount+", "Pepsi", "dentsu", "Stagwell"]
    for i, logo in enumerate(логотипы):
        row = i // 3; col = i % 3
        x = 0.70 + col * 2.20
        y = 2.65 + row * 0.65
        filled_rect(s, x, y, 2.05, 0.5, WHITE, stroke=LIGHT, stroke_pt=1.0,
                    radius=True, radius_adj=0.15)
        text_box(s, x=x, y=y, w=2.05, h=0.5,
                 text=logo, size=14, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Lionsgate × Runway quote with официальное объявление о партнёрстве image
    filled_rect(s, 0.70, 4.10, 6.7, 2.2, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.05)
    # Real Runway × Lionsgate баннер с объявлением on left
    add_image(s, str(ASSETS / "screenshots/s11-lionsgate-runway.png"),
              x=0.85, y=4.20, w=2.10, h=1.40)
    text_box(s, x=3.10, y=4.20, w=4.20, h=0.4,
             text="LIONSGATE × RUNWAY · SEP 18, 2024",
             size=10, bold=True, color=GOLD)
    text_box(s, x=3.10, y=4.65, w=4.20, h=1.4,
             text="\"...saving millions and millions of dollars on pre-production and post-production\"\n— Michael Burns, Vice Chairman of Lionsgate",
             size=11, italic=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=0.85, y=5.70, w=6.4, h=0.30,
             text="Source: investors.lionsgate.com · quarterly earnings call, Nov 2024",
             size=9, italic=True, color=SLATE)
    # Right: 3 metric chips
    text_box(s, x=7.85, y=2.0, w=5.1, h=0.4,
             text="METRICS 2026",
             size=11, bold=True, color=TEAL)
    metrics = [
        ("22B+", "assets in <2 years"),
        ("$400M", "direct revenue 2024-25"),
        ("40%", "of video ads — AI-generated (IAB 2026)"),
    ]
    for i, (val, label) in enumerate(metrics):
        y = 2.55 + i * 1.20
        filled_rect(s, 7.85, y, 5.1, 1.10, MID, radius=True, radius_adj=0.10)
        text_box(s, x=7.85, y=y + 0.10, w=5.1, h=0.55,
                 text=val, size=32, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=7.85, y=y + 0.70, w=5.1, h=0.35,
                 text=label, size=11, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Мини-провал
    lesson_box(s, 0.55, 6.6, 12.3, 0.85,
               "MINI-FAILURE: 86% of buyers use AI, but adoption ≠ success. Toys R Us Cannes 2024 — sentiment swing −10 pp.")
    speaker_notes(s, load_notes("s11"))


def build_s13(p):
    build_section_divider(p, here_idx=2, title="AI CHANGED",
                          frame_phrase="Cost collapse 100×-10,000× · speed days → seconds · new vs lost roles.",
                          notes_slide_id="s13")


def build_s14(p):
    """Обвал стоимости 100×-10 000×."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Cost collapse: 100×–10,000× across asset types",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="The marginal cost of generation is 2-4 orders of magnitude cheaper. But the commercially-safe enterprise segment is separate.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Cost table — 4 rows
    table_x = 0.55
    table_y = 2.0
    table_w = 8.5
    row_h = 0.90
    # Header
    filled_rect(s, table_x, table_y, table_w, 0.5, DEEP)
    headers = [("ASSET", 0, 2.5), ("BEFORE", 2.5, 2.0), ("AFTER", 4.5, 2.0), ("×", 6.5, 2.0)]
    for h, x_off, w in headers:
        text_box(s, x=table_x + x_off + 0.10, y=table_y + 0.10, w=w - 0.20, h=0.30,
                 text=h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("1 image", "$50-200 freelance\n$25-100 stock photo", "$0-0.25", "200×-10,000×"),
        ("50 product images", "$1-25K photoshoot", "$0-1.50", ">1,000×"),
        ("1 min of 720p video", "$1-50K shoot+post-production", "~$6 Sora 2", "150×-8,000×"),
        ("Dubbing $/min/language", "$50-500 actor+studio", "<$1 ElevenLabs", "50×-500×"),
    ]
    for i, (asset, before, after, mult) in enumerate(rows):
        y = table_y + 0.5 + i * row_h
        bg = SURFACE if i % 2 == 0 else WHITE
        filled_rect(s, table_x, y, table_w, row_h, bg,
                    stroke=LIGHT, stroke_pt=0.5)
        text_box(s, x=table_x + 0.10, y=y + 0.10, w=2.3, h=row_h - 0.20,
                 text=asset, size=12, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        text_box(s, x=table_x + 2.60, y=y + 0.10, w=1.8, h=row_h - 0.20,
                 text=before, size=11, italic=True, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        text_box(s, x=table_x + 4.60, y=y + 0.10, w=1.8, h=row_h - 0.20,
                 text=after, size=12, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        text_box(s, x=table_x + 6.60, y=y + 0.10, w=1.8, h=row_h - 0.20,
                 text=mult, size=13, bold=True, color=GOLD,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Right: $400M callout
    ocean_box(s, 9.25, 2.0, 3.65, 4.0, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.8)
    text_box(s, x=9.40, y=2.15, w=3.4, h=0.4,
             text="MID SEGMENT ≠ FREE",
             size=11, bold=True, color=GOLD)
    text_box(s, x=9.40, y=2.55, w=3.4, h=1.1,
             text="$400M",
             size=46, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=9.40, y=3.70, w=3.4, h=0.5,
             text="Adobe Firefly\ndirect revenue fiscal year 2024-25",
             size=11, color=DEEP, bold=True,
             align=PP_ALIGN.CENTER, line_spacing=1.30)
    text_box(s, x=9.40, y=4.50, w=3.4, h=1.4,
             text="licensed corpus + process + integration = enterprise SaaS stack\n\nThe lower segment is washed out, the mid segment grows.",
             size=10, italic=True, color=DEEP, line_spacing=1.30)
    # Footer
    text_box(s, x=0.55, y=6.75, w=12.3, h=0.4,
             text="Source: ZSky AI, Sora 2 API Pricing, ElevenLabs pricing.",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    """Speed-collapse: дни → секунды."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Speed: days → seconds",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="The iteration cycle becomes 10-100× tighter → new human skills.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 4 rows timer table
    rows = [
        ("Concept-art draft", "days (freelance designer)", "5-60 sec", "Midjourney / Flux / Imagen"),
        ("B-roll frame", "hours of shooting + post-production", "5-60 sec", "Veo / Sora"),
        ("Long-form dubbing", "weeks in a studio", "minutes", "ElevenLabs Dubbing Studio"),
        ("Concept exploration", "half a week · 3-5 options", "minutes · 10×+", "tighter iteration"),
    ]
    ocean_box(s, 0.55, 2.0, 12.3, 4.0)
    row_h = 0.85
    for i, (task, before, after, tool) in enumerate(rows):
        y = 2.20 + i * row_h
        # Task label
        text_box(s, x=0.80, y=y, w=2.8, h=row_h - 0.10,
                 text=task, size=14, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        # Before (Light teal — baseline, no red)
        filled_rect(s, 3.70, y + 0.10, 2.6, 0.55, SURFACE,
                    stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.20)
        text_box(s, x=3.70, y=y + 0.10, w=2.6, h=0.55,
                 text=before, size=11, italic=True, color=LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Arrow
        right_arrow(s, 6.45, y + 0.20, 0.6, 0.35, fill=MID)
        # After (gold accent — speed win)
        filled_rect(s, 7.20, y + 0.10, 2.0, 0.55, GOLD,
                    radius=True, radius_adj=0.20)
        text_box(s, x=7.20, y=y + 0.10, w=2.0, h=0.55,
                 text=after, size=15, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Tool
        text_box(s, x=9.35, y=y + 0.10, w=3.5, h=0.55,
                 text=tool, size=11, italic=True, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Anchor
    lesson_box(s, 0.55, 6.2, 12.3, 1.05,
               "Engineering lesson: the iteration cycle is 10-100× tighter → new human skills — deciding \"good enough / not good enough\" faster, because more new options appear.")
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    """New professions."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="New professions: prompt engineer · AI director · AI-workflow specialist",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="A specialized class between the AI tool and the final product (client deliverable).",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: Upwork screenshot мокап
    ocean_box(s, 0.55, 2.0, 5.5, 4.0)
    text_box(s, x=0.75, y=2.15, w=5.1, h=0.4,
             text="UPWORK · AI/ML CATEGORY",
             size=11, bold=True, color=TEAL)
    filled_rect(s, 0.75, 2.65, 5.1, 0.5, MID, radius=True, radius_adj=0.10)
    text_box(s, x=0.85, y=2.65, w=4.9, h=0.5,
             text="Search: AI / ML / prompt engineering",
             size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Mock results list
    jobs = [
        "Prompt engineer · Midjourney expert · $35-65/hr",
        "AI director / supervisor · $50-90/hr",
        "AI-workflow specialist · enterprise · $60-120/hr",
        "Continuity supervisor · video · $40-75/hr",
    ]
    for i, job in enumerate(jobs):
        y = 3.30 + i * 0.55
        filled_rect(s, 0.75, y, 5.1, 0.45, SURFACE,
                    stroke=LIGHT, stroke_pt=0.5, radius=True, radius_adj=0.18)
        text_box(s, x=0.85, y=y, w=4.9, h=0.45,
                 text=job, size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    # Metrics chip at bottom of card
    chip(s, 0.75, 5.55, 5.1, 0.40, "Upwork +70% year over year · 52% of services volume — AI",
         fill=GOLD_TINT, stroke=GOLD, color=DEEP, size=10, bold=True)
    # Right: 4 role-cards
    text_box(s, x=6.35, y=2.0, w=6.6, h=0.4,
             text="4 NEW ROLES",
             size=11, bold=True, color=TEAL)
    roles = [
        ("Prompt engineer / AI artist",
         "Shapes prompts + post-processing → a production-ready result",
         "$25-80/hr"),
        ("AI director / AI music producer",
         "Supervises model output: iterations + post-processing + multimodality",
         "Analogous to an art director"),
        ("AI-workflow specialist",
         "Integrator of AI tools into studios' production processes",
         "Adobe Firefly Foundry deployments"),
        ("Continuity supervisor",
         "Checks character/scene continuity across multi-frame sequences",
         "Controls drift between scenes"),
    ]
    for i, (rtitle, rbody, rextra) in enumerate(roles):
        y = 2.5 + i * 0.95
        ocean_box(s, 6.35, y, 6.6, 0.85)
        text_box(s, x=6.50, y=y + 0.05, w=6.3, h=0.30,
                 text=rtitle, size=12, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=6.50, y=y + 0.34, w=6.3, h=0.28,
                 text=rbody, size=10, color=SLATE, line_spacing=1.15)
        text_box(s, x=6.50, y=y + 0.60, w=6.3, h=0.22,
                 text=rextra, size=10, italic=True, color=GOLD, bold=True,
                 line_spacing=1.0)
    # Lesson box — repositioned to fit fully
    lesson_box(s, 0.55, 6.30, 12.3, 1.00,
               "New roles sit between the AI tool and the final product for the client. They grow fast, but are smaller than the displaced class.")
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """Замещение."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Displacement: −17% graphic design · Shutterstock isolated values",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Wage decline from the bottom is structural, not temporary. Industry consolidation as the response.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: big -17% callout (3 columns shrunk to make room for Урок full-width below)
    ocean_box(s, 0.55, 2.0, 4.0, 4.0)
    text_box(s, x=0.70, y=2.10, w=3.7, h=0.30,
             text="UPWORK · GRAPHIC DESIGN · YEAR OVER YEAR",
             size=10, bold=True, color=TEAL, line_spacing=1.0)
    text_box(s, x=0.70, y=2.50, w=3.7, h=1.30,
             text="−17%",
             size=68, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.70, y=3.85, w=3.7, h=0.30,
             text="of graphic-design jobs",
             size=11, italic=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=0.70, y=4.25, w=3.7, h=1.65,
             text="40% of copywriter jobs $10-19/hr — detected as AI · <10% at $60+/hr · wage decline from the bottom",
             size=10, color=SLATE, line_spacing=1.30, align=PP_ALIGN.CENTER)
    # Middle: Shutterstock licensing timeline
    ocean_box(s, 4.75, 2.0, 4.0, 4.0)
    text_box(s, x=4.90, y=2.10, w=3.7, h=0.30,
             text="SHUTTERSTOCK · LICENSING FOR AI",
             size=10, bold=True, color=TEAL)
    # Timeline bars
    years_data = [("2023", 104, MID), ("2024", 138, LIGHT), ("2027", 250, GOLD)]
    max_val = 250
    for i, (year, val, color) in enumerate(years_data):
        y = 2.55 + i * 0.85
        bar_w = (val / max_val) * 3.5
        text_box(s, x=4.90, y=y, w=0.6, h=0.35,
                 text=year, size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        filled_rect(s, 5.55, y, bar_w, 0.45, color, radius=True, radius_adj=0.20)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=5.65, y=y, w=bar_w - 0.1, h=0.45,
                 text=f"${val}M", size=13, bold=True, color=text_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=4.90, y=5.45, w=3.7, h=0.30,
             text="pivot: photos → training data for AI",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right: Getty merger + SAG-AFTRA
    ocean_box(s, 8.95, 2.0, 4.0, 4.0)
    text_box(s, x=9.10, y=2.10, w=3.7, h=0.30,
             text="INDUSTRY RESPONSE",
             size=10, bold=True, color=TEAL)
    # Getty/Shutterstock merger
    filled_rect(s, 9.10, 2.50, 3.7, 0.90, MID, radius=True, radius_adj=0.10)
    text_box(s, x=9.20, y=2.55, w=3.5, h=0.35,
             text="Getty + Shutterstock merger", size=10, bold=True, color=WHITE)
    text_box(s, x=9.20, y=2.92, w=3.5, h=0.45,
             text="$3.7B · Jan 2025 · Defensive consolidation",
             size=9, italic=True, color=WHITE, line_spacing=1.20)
    # SAG-AFTRA chip (Screen Actors Guild) + WGA (Writers Guild)
    filled_rect(s, 9.10, 3.50, 3.7, 0.90, LIGHT, radius=True, radius_adj=0.10)
    text_box(s, x=9.20, y=3.55, w=3.5, h=0.35,
             text="SAG-AFTRA + WGA*", size=10, bold=True, color=WHITE)
    text_box(s, x=9.20, y=3.92, w=3.5, h=0.45,
             text="Digital-replica clause · 2026 · 4-year extension",
             size=9, italic=True, color=WHITE, line_spacing=1.20)
    # Voice actors
    filled_rect(s, 9.10, 4.50, 3.7, 0.90, GOLD, radius=True, radius_adj=0.10)
    text_box(s, x=9.20, y=4.55, w=3.5, h=0.35,
             text="Voice actors · worldwide", size=10, bold=True, color=DEEP)
    text_box(s, x=9.20, y=4.92, w=3.5, h=0.45,
             text="Commodity dubbing displaced · ElevenLabs in the corporate sector",
             size=9, italic=True, color=DEEP, line_spacing=1.20)
    # Glossary footer (smaller, bottom of right column)
    text_box(s, x=9.10, y=5.55, w=3.7, h=0.30,
             text="* SAG-AFTRA = Screen Actors Guild · WGA = Writers Guild of America",
             size=8, italic=True, color=LIGHT, align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Lesson — full width below all 3 columns
    lesson_box(s, 0.55, 6.20, 12.3, 1.10,
               "Structurally: AI washes out the lower segment and leaves the top segment protected. The clauses do not cover bottom-tier freelance.")
    speaker_notes(s, load_notes("s17"))


def build_s19(p):
    build_section_divider(p, here_idx=3, title="AI BROKE",
                          frame_phrase="12 cases: copyright × 4 + deepfakes × 2 + slop · fake authors · backlash · displacement + taxonomy.",
                          notes_slide_id="s19")


def build_s20(p):
    """Copyright 4 categories таксономия."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="\"AI and copyright\" — 4 different categories of lawsuit",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Not one legal question, but four. With different legal logic and outcomes.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 2x2 matrix
    grid_x = 0.55
    grid_y = 2.0
    grid_w = 8.5
    grid_h = 4.3
    cell_w = (grid_w - 0.20) / 2
    cell_h = (grid_h - 0.20) / 2
    cats = [
        ("1. Scraping for training", "INPUT SIDE",
         "An AI company assembled a corpus WITHOUT a license",
         "→ NYT v. OpenAI · Andersen", MID),
        ("2. Output similarity", "OUTPUT SIDE",
         "The model reproduces protected content verbatim",
         "→ verbatim quotation of NYT · DMCA", LIGHT),
        ("3. Style imitation", "\"IN THE STYLE OF A SPECIFIC ARTIST\"",
         "Class action by artists · style is not copyrightable, but DMCA + rights of publicity",
         "→ Andersen v. Stability/MJ/Deviant", GOLD),
        ("4. Voice/likeness", "RIGHT OF PUBLICITY",
         "Use of a voice/likeness without consent",
         "→ ScarJo v. OpenAI · SAG-AFTRA · Korea", TEAL),
    ]
    for i, (title, subtitle, body, cases, accent) in enumerate(cats):
        row = i // 2; col = i % 2
        x = grid_x + col * (cell_w + 0.20)
        y = grid_y + row * (cell_h + 0.20)
        ocean_box(s, x, y, cell_w, cell_h)
        # Accent corner
        filled_rect(s, x, y, cell_w, 0.55, accent, radius=True, radius_adj=0.10)
        accent_text = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=y + 0.08, w=cell_w - 0.30, h=0.40,
                 text=title, size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.18, y=y + 0.65, w=cell_w - 0.36, h=0.35,
                 text=subtitle, size=10, bold=True, italic=True, color=TEAL)
        text_box(s, x=x + 0.18, y=y + 1.00, w=cell_w - 0.36, h=0.7,
                 text=body, size=12, color=DEEP, line_spacing=1.30)
        text_box(s, x=x + 0.18, y=y + 1.65, w=cell_w - 0.36, h=0.4,
                 text=cases, size=10, italic=True, color=GOLD, bold=True,
                 line_spacing=1.20)
    # Right: lesson box
    lesson_box(s, 9.30, 2.0, 3.65, 4.3,
               "\"AI copyright\" is NOT one question, but 4 different categories of risk.\n\nCheck which of the 4 applies to your process.")
    # Footer
    text_box(s, x=0.55, y=6.55, w=12.3, h=0.35,
             text="Each of the 4 categories is examined next through a separate reference case.",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s20"))


def case_slide_template(p, slide_id, title, assertion_body, screenshot_path,
                        timeline_events, lesson_text, screenshot_label=None,
                        emphasis_block=None, glossary=None):
    """Generic case-slide builder for s21-s25, s27. Phase 6+7 v2.

    screenshot_path — abs path to real news-card PNG (REQUIRED).
    emphasis_block — optional dict for slide-specific emphasis layout:
        {"kind": "big_number", "value": "20M", "label": "ChatGPT logs"}
        {"kind": "verdict_badge", "value": "STABILITY ВЫИГРАЛ", "caption": "..."}
        {"kind": "trial_chip", "value": "8 СЕНТ 2026", "label": "ДАТА СУДА"}
    glossary — short инлайн-пояснение (e.g. «MTD = motion to dismiss»)
    """
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text=title, size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text=assertion_body, size=13, italic=True, color=MID, align=PP_ALIGN.LEFT,
             line_spacing=1.30)
    # Left: real news-screenshot embedded
    ocean_box(s, 0.55, 2.0, 6.0, 4.0)
    if screenshot_label:
        text_box(s, x=0.70, y=2.05, w=5.7, h=0.3,
                 text=screenshot_label, size=10, bold=True, color=TEAL)
    if screenshot_path and Path(screenshot_path).exists():
        add_image(s, screenshot_path, x=0.70, y=2.40, w=5.70, h=3.50)
    # Right: timeline events
    ocean_box(s, 6.85, 2.0, 6.10, 4.0)
    text_box(s, x=7.00, y=2.10, w=5.8, h=0.4,
             text="TIMELINE", size=11, bold=True, color=TEAL)
    for i, (date, event, color) in enumerate(timeline_events):
        y = 2.60 + i * 0.78
        # Date chip
        filled_rect(s, 7.00, y, 1.5, 0.5, color, radius=True, radius_adj=0.20)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=7.00, y=y, w=1.5, h=0.5, text=date,
                 size=10, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Event text
        text_box(s, x=8.60, y=y, w=4.3, h=0.5, text=event,
                 size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Emphasis block (slide-specific diversification)
    if emphasis_block:
        eb = emphasis_block
        if eb.get("kind") == "big_number":
            # Big number callout below timeline
            filled_rect(s, 6.85, 5.65, 6.10, 0.45, GOLD,
                        radius=True, radius_adj=0.20)
            text_box(s, x=6.85, y=5.65, w=6.10, h=0.45,
                     text=f"{eb['value']}  ·  {eb['label']}",
                     size=12, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "verdict_badge":
            filled_rect(s, 6.85, 5.65, 6.10, 0.45, TEAL,
                        radius=True, radius_adj=0.20)
            text_box(s, x=6.85, y=5.65, w=6.10, h=0.45,
                     text=f"{eb['value']}",
                     size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "trial_chip":
            filled_rect(s, 6.85, 5.65, 6.10, 0.45, MID,
                        radius=True, radius_adj=0.20)
            text_box(s, x=6.85, y=5.65, w=6.10, h=0.45,
                     text=f"{eb['label']} · {eb['value']}",
                     size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "settlement_matrix":
            # 3-major × 2-defendant matrix (2 rows of 3, или 1 row of 3 для legacy)
            label_y = 5.40
            filled_rect(s, 6.85, label_y, 6.10, 0.28, SURFACE,
                        radius=True, radius_adj=0.30)
            text_box(s, x=6.85, y=label_y, w=6.10, h=0.28,
                     text=eb.get("title", "MAJOR LABELS · SETTLEMENT STATUS"),
                     size=9, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cells = eb.get("cells", [])
            # Layout: 3 columns x N rows
            ncols = 3
            nrows = (len(cells) + ncols - 1) // ncols
            cw = 6.10 / ncols
            ch = 0.32
            for ci, (label, color) in enumerate(cells):
                row = ci // ncols
                col = ci % ncols
                cx = 6.85 + col * cw
                cy = label_y + 0.32 + row * (ch + 0.05)
                filled_rect(s, cx + 0.05, cy, cw - 0.10, ch,
                            color, radius=True, radius_adj=0.30)
                ct = WHITE if color != GOLD else DEEP
                text_box(s, x=cx, y=cy, w=cw, h=ch,
                         text=label, size=8, bold=True, color=ct,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "fair_use_factors":
            label_y = 5.50
            filled_rect(s, 6.85, label_y, 6.10, 0.30, SURFACE,
                        radius=True, radius_adj=0.30)
            text_box(s, x=6.85, y=label_y, w=6.10, h=0.30,
                     text="\"FAIR USE\" 4-FACTOR TEST (Warhol v. Goldsmith)",
                     size=9, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            factors = eb.get("factors", [])
            fw = 6.10 / max(1, len(factors))
            for fi, (label, status_color) in enumerate(factors):
                fx = 6.85 + fi * fw
                filled_rect(s, fx + 0.05, label_y + 0.35, fw - 0.10, 0.30,
                            status_color, radius=True, radius_adj=0.30)
                ct = WHITE if status_color != GOLD else DEEP
                text_box(s, x=fx, y=label_y + 0.35, w=fw, h=0.30,
                         text=label, size=9, bold=True, color=ct,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Optional glossary footnote
    if glossary:
        text_box(s, x=6.85, y=6.05, w=6.10, h=0.25,
                 text=glossary, size=8, italic=True, color=LIGHT,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Lesson box at bottom
    lesson_box(s, 0.55, 6.30, 12.3, 1.05, lesson_text)
    speaker_notes(s, load_notes(slide_id))


def build_s21(p):
    # Emphasis: 20M ChatGPT logs (big number callout)
    case_slide_template(p, "s21",
        title="NYT v. OpenAI — training + output similarity (Case 1)",
        assertion_body="NYT v. OpenAI (Dec 2023): the court ordered OpenAI to produce 20M ChatGPT logs. Summary judgment (SJ*) deadline — Apr 2, 2026. Verbatim-quotation theory.",
        screenshot_path=str(ASSETS / "screenshots/s21-nyt-bloomberg.png"),
        screenshot_label="THE NEW YORK TIMES · DEC 27, 2023",
        timeline_events=[
            ("Dec 2023", "NYT filed suit", MID),
            ("2024-25", "Discovery + fight over procedural motions", LIGHT),
            ("Apr 2, 2026", "Summary judgment (SJ) deadline", GOLD),
        ],
        emphasis_block={"kind": "big_number", "value": "20,000,000", "label": "ChatGPT logs under discovery"},
        glossary="* SJ = summary judgment (a ruling without a full trial)",
        lesson_text="If the model can quote your training corpus verbatim — that is NOT \"fair use\", it is evidence of infringement. Checking output similarity against the training data is mandatory."
    )


def build_s22(p):
    # Emphasis: UK verdict badge "STABILITY WON"
    case_slide_template(p, "s22",
        title="Getty v. Stability — a UK win vs a US wait (Case 2)",
        assertion_body="UK High Court 11/04/2025 — Stability won the main claims (weights ≠ a copy under CDPA*). The US case — motion to dismiss** 02/10/2026.",
        screenshot_path=str(ASSETS / "screenshots/s22-getty-bird.png"),
        screenshot_label="THE VERGE · GETTY vs STABILITY · FEB 2023",
        timeline_events=[
            ("Nov 4, 2025", "UK · Stability won the main claims", TEAL),
            ("Feb 10, 2026", "US case · motion to dismiss expected", GOLD),
            ("Pending", "Trademark + passing-off — separate claims", LIGHT),
        ],
        emphasis_block={"kind": "verdict_badge",
                        "value": "UK: STABILITY WON the main claims under CDPA"},
        glossary="* CDPA = UK Copyright, Designs and Patents Act 1988  ·  ** MTD = motion to dismiss",
        lesson_text="Jurisdictions diverge — what is legal in the UK under CDPA is ILLEGAL in the US under \"fair use\". For a global deployment, check both."
    )


def build_s23(p):
    # Emphasis: trial date prominent
    case_slide_template(p, "s23",
        title="Andersen v. Stability/MJ/Deviant — style imitation (Case 3)",
        assertion_body="Class action by artists. The motion to dismiss* was denied in Aug 2024 → discovery. Third amended complaint — Feb 27, 2026. Trial — Sep 8, 2026.",
        screenshot_path=str(ASSETS / "screenshots/s23-andersen-docket.png"),
        screenshot_label="KELLY McKERNAN (PLAINTIFF) · WIKIMEDIA COMMONS",
        timeline_events=[
            ("Jan 2023", "Class action filed (10 artists)", MID),
            ("Aug 2024", "motion to dismiss denied (Judge Orrick) → discovery", LIGHT),
            ("Feb 27, 2026", "3rd amended complaint", TEAL),
        ],
        emphasis_block={"kind": "trial_chip", "label": "TRIAL DATE", "value": "SEP 8, 2026"},
        glossary="* MTD = motion to dismiss",
        lesson_text="Style imitation \"in the style of [a specific artist]\" is UNSAFE, even if style is not copyrightable. Class actions survive the motion to dismiss on DMCA + rights of publicity."
    )


def build_s24(p):
    # Emphasis: 3 major × 2 defendant matrix
    case_slide_template(p, "s24",
        title="RIAA v. Suno/Udio — licensing under litigation pressure (Case 4)",
        assertion_body="RIAA filed suit 06/24/2024. UMG settled with Udio 10/29/2025; Warner settled with Suno Sep 2025. Sony is actively litigating against both.",
        screenshot_path=str(ASSETS / "screenshots/s24-riaa-suno.png"),
        screenshot_label="BILLBOARD · MAJOR LABEL LAWSUIT · JUN 24, 2024",
        timeline_events=[
            ("Jun 24, 2024", "RIAA files suit against Suno + Udio", MID),
            ("Oct 29, 2025", "UMG × Udio settlement → joint platform", TEAL),
            ("Sep 2025", "Warner × Suno settlement (royalties + stake)", TEAL),
            ("Jul 2026", "Suno SJ — Sony actively litigating against both", GOLD),
        ],
        emphasis_block={"kind": "settlement_matrix",
                        "title": "3 MAJOR LABELS × 2 DEFENDANTS — status",
                        "cells": [
                            ("UMG × Udio: settled", TEAL),
                            ("UMG × Suno: negotiating", LIGHT),
                            ("Warner × Suno: settled", TEAL),
                            ("Warner × Udio: litigating", GOLD),
                            ("Sony × Suno: litigating", GOLD),
                            ("Sony × Udio: litigating", GOLD),
                        ]},
        glossary="* UMG = Universal Music Group (one of the 3 \"Big Three\" major labels)",
        lesson_text="Licensing under litigation pressure — the actual outcome: 4 of 6 plaintiff-defendant combinations are settled or in negotiation. This is a new layer of the business model, not a \"ban on all AI music\"."
    )


def build_s25(p):
    # Emphasis: «добросовестное использование» — 4-факторный тест breakdown
    case_slide_template(p, "s25",
        title="Thomson Reuters v. Ross — the first US denial of \"fair use\" (Case 5)",
        assertion_body="Feb 2025, Judge Bibas: 2200/3000 headnotes — infringement, \"fair use\" rejected. Caveat: Ross is a non-generative AI.",
        screenshot_path=str(ASSETS / "screenshots/s25-thomson-reedsmith.png"),
        screenshot_label="DAVIS WRIGHT TREMAINE · ROSS RULING · FEB 2025",
        timeline_events=[
            ("Feb 2025", "Bibas · the first US denial of \"fair use\"", GOLD),
            ("2200/3000", "headnotes — infringement", LIGHT),
            ("Pending", "Test cases for LLM/diffusion still ahead", MID),
        ],
        emphasis_block={"kind": "fair_use_factors",
                        "factors": [
                            ("1 · purpose", GOLD),
                            ("2 · nature", GOLD),
                            ("3 · amount", GOLD),
                            ("4 · market", GOLD),
                        ]},
        glossary="Caveat: Ross — non-generative AI (legal search). Generative LLM/diffusion test cases are still ahead.",
        lesson_text="\"Fair use\" is NOT the default. Test cases for LLM/diffusion are still ahead. Do NOT build a product roadmap on the assumption that \"fair use\" will work as a defense."
    )


def build_s26(p):
    """Arup deepfake — special schema (attack diagram). NOTE: lesson_box already has УРОК prefix so we strip it from text."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Arup CFO deepfake — a $25.6M fraud (Case 6)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Hong Kong, Jan 2024 · A finance worker on a video call with a deepfake CFO + colleagues → 15 transactions.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: CNN главное изображение — actual published photo (hands on laptop / hacker imagery)
    ocean_box(s, 0.55, 2.0, 4.5, 4.0)
    text_box(s, x=0.70, y=2.08, w=4.2, h=0.30,
             text="CNN · $25M FRAUD IN HONG KONG · MAY 16, 2024", size=10, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s26-arup-cnn.png"),
              x=0.70, y=2.40, w=4.20, h=3.30)
    text_box(s, x=0.70, y=5.75, w=4.2, h=0.30,
             text="Arup · British engineering firm · Sydney Opera House",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right: attack diagram — 5 stages (Gold replaces Red on final step)
    text_box(s, x=5.35, y=2.0, w=7.6, h=0.4,
             text="ATTACK SCENARIO", size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    stages = [
        ("1", "Email from the \"CFO\"", MID),
        ("2", "Invitation to a video call", LIGHT),
        ("3", "Call with deepfakes", TEAL),
        ("4", "15 transactions", MID),
        ("5", "$25.6M lost", GOLD),
    ]
    stage_y = 2.55
    stage_w = 1.30
    arrow_w = 0.20
    total_units = len(stages) * stage_w + (len(stages) - 1) * arrow_w
    start_x = 5.35 + (7.6 - total_units) / 2
    for i, (num, label, color) in enumerate(stages):
        x = start_x + i * (stage_w + arrow_w)
        filled_rect(s, x, stage_y, stage_w, 2.5, color, radius=True, radius_adj=0.10)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=x, y=stage_y + 0.20, w=stage_w, h=0.6, text=num,
                 size=32, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(s, x=x + 0.05, y=stage_y + 1.0, w=stage_w - 0.10, h=1.4,
                 text=label, size=10, color=text_color, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)
        if i < len(stages) - 1:
            right_arrow(s, x + stage_w + 0.02, stage_y + 1.0,
                        arrow_w - 0.04, 0.5, fill=DEEP)
    # Lesson
    lesson_box(s, 0.55, 6.4, 12.3, 0.95,
               "A video call ≠ identity verification in 2024+. Financial transactions require verification through an independent channel — a callback to a known number, multi-factor authentication, a documented process.")
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """Korea deepfake crisis — text-only, NO deepfake visuals."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="The schoolgirl deepfake crisis in Korea (Case 7)",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Aug 2024 · >200 Telegram chats from selfies of classmates/teachers · 4× vs 2023.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: PBS NewsHour real photo — masked protestors with «반복되는 딥페이크 성범죄 국가도 공범이다»
    # banner (=«Repeated deepfake sex crimes, the state is an accomplice»). NO deepfake visuals.
    ocean_box(s, 0.55, 2.0, 5.0, 4.0)
    text_box(s, x=0.70, y=2.08, w=4.7, h=0.30,
             text="PBS NEWSHOUR · PROTEST IN KOREA · 2024", size=10, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s27-korea-npr.png"),
              x=0.70, y=2.40, w=4.70, h=3.50)
    # Right: 4 data cards
    text_box(s, x=5.85, y=2.0, w=7.1, h=0.4,
             text="CRISIS FIGURES", size=11, bold=True, color=TEAL)
    stats = [
        (">200", "Telegram chats with deepfake porn"),
        ("6,500", "takedown requests, Jan-Jul 2024 (4× vs 2023)"),
        ("74%", "of suspects — 10-19 years old"),
        ("793 / 16", "reports / criminal cases (2021 — Jul 2024)"),
    ]
    for i, (val, label) in enumerate(stats):
        y = 2.55 + i * 0.95
        ocean_box(s, 5.85, y, 7.1, 0.85)
        text_box(s, x=6.00, y=y + 0.10, w=2.0, h=0.65,
                 text=val, size=28, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=8.10, y=y + 0.10, w=4.7, h=0.65,
                 text=label, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.30)
    # Lesson
    lesson_box(s, 0.55, 6.4, 12.3, 0.95,
               "An accessible capability + weak enforcement = mass collective harm. For consumer AI tools, a safety layer (NSFW-content detection + age verification + complaint-handling process) is mandatory BEFORE launch.")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """Slop + коллапс моделей."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Slop + model collapse · Google AI Overviews (Case 8)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="\"Put glue on pizza\" + \"eat one rock per day\". Shumailov, Nature 2024: recursive training → degradation.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: 2 Google AI Overview screenshots — Ocean palette
    ocean_box(s, 0.55, 2.0, 7.0, 4.2)
    text_box(s, x=0.70, y=2.10, w=6.7, h=0.4,
             text="GOOGLE AI OVERVIEW · real answers (May 2024)",
             size=10, bold=True, color=TEAL)
    # Screenshot 1
    filled_rect(s, 0.70, 2.55, 3.30, 1.7, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.06)
    text_box(s, x=0.85, y=2.65, w=3.10, h=0.40,
             text="\"how to keep cheese\nfrom sliding off pizza\"",
             size=10, bold=True, color=DEEP, line_spacing=1.20)
    # Gold accent (anti-pattern flag) — not red
    filled_rect(s, 0.78, 3.18, 0.04, 0.95, fill=GOLD)
    text_box(s, x=0.95, y=3.20, w=3.00, h=1.0,
             text="→ AI Overview: \"add ⅛ cup of non-toxic glue to the sauce\"",
             size=10, italic=True, color=DEEP, bold=True, line_spacing=1.30)
    text_box(s, x=0.85, y=4.05, w=3.10, h=0.20,
             text="source: a Reddit joke (11 years old)",
             size=9, italic=True, color=LIGHT)
    # Screenshot 2
    filled_rect(s, 4.20, 2.55, 3.20, 1.7, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.06)
    text_box(s, x=4.35, y=2.65, w=2.9, h=0.40,
             text="\"how many rocks should\nI eat per day?\"",
             size=10, bold=True, color=DEEP, line_spacing=1.20)
    filled_rect(s, 4.28, 3.18, 0.04, 0.95, fill=GOLD)
    text_box(s, x=4.45, y=3.20, w=2.80, h=1.0,
             text="→ AI Overview: \"at least one small rock per day\"",
             size=10, italic=True, color=DEEP, bold=True, line_spacing=1.30)
    text_box(s, x=4.35, y=4.05, w=2.9, h=0.20,
             text="source: The Onion (satire)",
             size=9, italic=True, color=LIGHT)
    # Below: Nature paper card with real Fig 1 (perplexity histograms коллапс моделей)
    filled_rect(s, 0.70, 4.50, 6.7, 1.55, MID, radius=True, radius_adj=0.08)
    # Real Nature paper Fig 1 — Shumailov 2024 коллапс моделей figure on left
    add_image(s, str(ASSETS / "screenshots/s28-ai-overview-glue.png"),
              x=0.80, y=4.60, w=1.40, h=1.35)
    text_box(s, x=2.30, y=4.60, w=5.0, h=0.40,
             text="NATURE · vol 631 · p 755-759 (2024)",
             size=11, bold=True, color=GOLD)
    text_box(s, x=2.30, y=5.05, w=5.0, h=0.55,
             text="Shumailov et al. — \"AI models collapse when trained on recursively generated data\"",
             size=11, color=WHITE, bold=True, line_spacing=1.25)
    text_box(s, x=0.85, y=5.65, w=6.4, h=0.35,
             text="MAD: Model Autophagy Disorder",
             size=11, italic=True, color=WHITE, line_spacing=1.0)
    # Right: lesson
    lesson_box(s, 7.70, 2.0, 5.25, 4.2,
               "Source quality matters more than volume.\n\nA model trained on unfiltered Reddit jokes loses to a model trained on a curated dataset — even if the curated one is 10× smaller.\n\nThis explains why Adobe Firefly works: a licensed corpus, not data scraped from the web.")
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """SI — фейковые авторы + Amazon."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Sports Illustrated fake authors + Amazon fake books (Case 9)",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="SI, Nov 2023 · articles with AI-generated profile photos. Authors Guild: a surge of fake books on Amazon Kindle under AI pseudonyms.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: CNN article screenshot — actual Drew Ortiz фейковый профиль from SI website
    # (Sports Illustrated bio page with AI face + AI-сгенерированный outdoors-reviewer text)
    ocean_box(s, 0.55, 2.0, 5.5, 4.0)
    text_box(s, x=0.70, y=2.08, w=5.2, h=0.30,
             text="CNN · SPORTS ILLUSTRATED FAKE AUTHORS · NOV 27, 2023", size=10, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s29-si-futurism.png"),
              x=0.70, y=2.40, w=5.20, h=2.65)
    chip(s, 0.70, 5.20, 5.2, 0.4,
         "\"Drew Ortiz\" — entirely fictional · AI portrait + AI text",
         fill=GOLD, color=DEEP, size=10, bold=True)
    # Right: Amazon Kindle Authors Guild data (no fabricated 19/100 number)
    ocean_box(s, 6.30, 2.0, 6.65, 4.0)
    text_box(s, x=6.45, y=2.10, w=6.3, h=0.4,
             text="AMAZON KINDLE 2023-24 · AUTHORS GUILD",
             size=11, bold=True, color=TEAL)
    # Big icon/headline replaces fabricated 19/100 stat
    text_box(s, x=6.45, y=2.55, w=6.3, h=1.0,
             text="AI PSEUDONYMS",
             size=44, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=6.45, y=3.55, w=6.3, h=0.5,
             text="a surge of fake books (Authors Guild 2023-24)",
             size=13, italic=True, color=DEEP, align=PP_ALIGN.CENTER, bold=True)
    text_box(s, x=6.45, y=4.10, w=6.3, h=1.85,
             text="AI-generated knockoffs pass themselves off as real jazz figures, financial advisors\n\n\"Frank Gioia\" · \"Ted Alkyer\" — fakes of real jazz personalities\n\nAmazon imposed a limit of 3 books/day + AI disclosure",
             size=11, color=SLATE, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.40)
    # Lesson
    lesson_box(s, 0.55, 6.20, 12.3, 1.0,
               "Accumulated trust is a key brand asset. AI pseudonyms destroy it instantly. If you publish under a name — the name must be a real person OR the AI authorship must be explicitly disclosed.")
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """Toys R Us / Coca-Cola разворот тональности — real ad screenshots + delta."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Marketing failure — Toys R Us + Coca-Cola (Case 10)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Toys R Us Cannes Lions 2024 — sentiment swing −10 pp. Coca-Cola Holidays 2024 AI ad — viral negativity.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: Toys R Us official Sora ad still
    ocean_box(s, 0.55, 2.0, 6.0, 2.55)
    text_box(s, x=0.70, y=2.10, w=5.7, h=0.4,
             text="TOYS R US STUDIOS · SORA AD · CANNES 2024",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s30-toysrus-cannes.png"),
              x=0.70, y=2.45, w=5.70, h=2.00)
    # Right: Coca-Cola Holidays Are Coming AI-реклама still
    ocean_box(s, 6.85, 2.0, 6.0, 2.55)
    text_box(s, x=7.00, y=2.10, w=5.7, h=0.4,
             text="COCA-COLA · «HOLIDAYS ARE COMING» AI · 2024",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s30-coca-secret.png"),
              x=7.00, y=2.45, w=5.70, h=2.00)
    # Two разворот тональности chips below
    # Toys R Us swing
    filled_rect(s, 0.55, 4.70, 6.0, 1.60, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.08)
    text_box(s, x=0.70, y=4.78, w=5.7, h=0.30,
             text="TOYS R US · SENTIMENT SWING (June 2024)",
             size=10, bold=True, color=TEAL)
    text_box(s, x=0.70, y=5.10, w=2.9, h=0.5,
             text="POSITIVE", size=10, color=SLATE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=0.70, y=5.45, w=2.9, h=0.5,
             text="+12.2%  →  +3.4%", size=14, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=3.65, y=5.10, w=2.7, h=0.5,
             text="NEGATIVE", size=10, color=SLATE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=3.65, y=5.45, w=2.7, h=0.5,
             text="13.5%  →  53.4%", size=14, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=0.70, y=5.90, w=5.7, h=0.30,
             text="−8.8 pp positive · +39.9 pp negative", size=10, italic=True, color=DEEP)
    # Coca-Cola swing
    filled_rect(s, 6.85, 4.70, 6.0, 1.60, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.08)
    text_box(s, x=7.00, y=4.78, w=5.7, h=0.30,
             text="COCA-COLA AI AD · NEGATIVE (Nov 2024)",
             size=10, bold=True, color=TEAL)
    text_box(s, x=7.00, y=5.10, w=5.7, h=0.5,
             text="\"soulless\" · \"creepy faces\" · \"the truck wheels spin the wrong way\"",
             size=11, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    text_box(s, x=7.00, y=5.65, w=5.7, h=0.5,
             text="100 employees · 70,000 AI clips · viral mockery",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    # Lesson
    lesson_box(s, 0.55, 6.5, 12.3, 0.9,
               "AI advertising is possible, but a flagship seasonal campaign WITHOUT human leadership = brand damage. Sentiment swing, NOT CTR, is the primary metric of brand-trust risk.")
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """Замещение consolidated 3-stat."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Displacement: the consolidated picture — a structural shift (Case 11)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Wage decline from the bottom is structural, not a temporary shock. The clauses help the top segment, not the bottom.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 stat blocks horizontal
    stats = [
        ("−17.01%", "graphic-design jobs\nUpwork year over year", MID, "Jobbers Index"),
        ("40% / <10%", "jobs $10-19/hr vs\n$60+/hr · detected as AI",
         LIGHT, "wage decline from the bottom"),
        ("4 years", "SAG-AFTRA + WGA\nextension 2026", GOLD,
         "Digital-replica clause"),
    ]
    card_y = 2.2
    card_h = 3.6
    card_w = 3.95
    gap = 0.18
    start_x = 0.55
    for i, (big, label, color, footer) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        # Color band
        filled_rect(s, x, card_y, card_w, 0.6, color, radius=True, radius_adj=0.10)
        # Big stat
        text_box(s, x=x + 0.15, y=card_y + 0.8, w=card_w - 0.30, h=1.4,
                 text=big, size=44, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Label
        text_box(s, x=x + 0.15, y=card_y + 2.3, w=card_w - 0.30, h=0.7,
                 text=label, size=13, color=DEEP, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.30)
        # Footer chip
        chip(s, x + 0.30, card_y + card_h - 0.55, card_w - 0.60, 0.40, footer,
             fill=SURFACE, stroke=LIGHT, color=DEEP, size=10, bold=False)
    # Getty merger chip below
    text_box(s, x=0.55, y=6.05, w=12.3, h=0.4,
             text="+ Getty + Shutterstock merger $3.7B (Jan 2025) — defensive consolidation",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Lesson
    lesson_box(s, 0.55, 6.45, 12.3, 0.75,
               "Displacement is structural, not a temporary shock. The clauses help, but wage decline from the bottom remains. Understand which class of workers your AI project will displace BEFORE launch.")
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    build_section_divider(p, here_idx=4, title="AI not needed here",
                          frame_phrase="Criteria for negative selection · 3 human-only zones · empirical rejection by end users.",
                          notes_slide_id="s32")


def build_s33(p):
    """4 критерия отказа от AI."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="4 criteria for refusing AI in a creative project",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="The lessons learned from Section 3 → translated into a checklist.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 2x2 grid of criteria
    grid_x = 0.55
    grid_y = 2.0
    grid_w = 12.3
    grid_h = 4.6
    cell_w = (grid_w - 0.20) / 2
    cell_h = (grid_h - 0.20) / 2
    criteria = [
        ("1. License on training data",
         "No documented licensed corpus → legal debt under category 1.",
         "→ Andersen v. Stability · RIAA v. Suno",
         "Firefly = yes · Stable Diffusion = risks",
         MID),
        ("2. Output-similarity check",
         "Outputs may reproduce protected content → liability under category 2.",
         "→ NYT verbatim-quotation theory",
         "A technical control is mandatory",
         LIGHT),
        ("3. Consent for voice/likeness",
         "No explicit consent → risk of the ScarJo + SAG-AFTRA + Korea class.",
         "→ ScarJo v. OpenAI · the Korea crisis",
         "Recognizable real people — always consent",
         TEAL),
        ("4. Brand-trust risk",
         "A flagship/historical campaign without human leadership → measurable backlash.",
         "→ Coca-Cola · Toys R Us · SI",
         "Sentiment swing, not CTR",
         GOLD),
    ]
    for i, (title, body, cases, tip, accent) in enumerate(criteria):
        row = i // 2; col = i % 2
        x = grid_x + col * (cell_w + 0.20)
        y = grid_y + row * (cell_h + 0.20)
        ocean_box(s, x, y, cell_w, cell_h)
        filled_rect(s, x, y, cell_w, 0.55, accent, radius=True, radius_adj=0.10)
        accent_text = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=y + 0.08, w=cell_w - 0.30, h=0.40,
                 text=title, size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.18, y=y + 0.70, w=cell_w - 0.36, h=0.60,
                 text=body, size=11, color=DEEP, line_spacing=1.25)
        text_box(s, x=x + 0.18, y=y + 1.40, w=cell_w - 0.36, h=0.35,
                 text=cases, size=10, italic=True, color=GOLD, bold=True,
                 line_spacing=1.15)
        # Tip chip — separate solid block to avoid overlap
        chip(s, x + 0.18, y + cell_h - 0.50, cell_w - 0.36, 0.38, tip,
             fill=TEAL_TINT, stroke=TEAL, color=DEEP, size=10, bold=True)
    speaker_notes(s, load_notes("s33"))


def build_s34(p):
    """3 zones только человеком."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="3 zones where AI should not replace the human",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="\"AI as a tool\" works; \"AI as replacement\" does not.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 columns
    card_y = 2.0
    card_h = 4.6
    card_w = (12.3 - 0.30) / 3
    start_x = 0.55
    gap = 0.15
    zones = [
        ("Investigative journalism",
         "NYT/WaPo guidelines prohibit AI for primary reporting",
         [
             "Source verification — human-only",
             "On-the-record interviews — human-only",
             "Accountability journalism — human responsibility is non-delegable",
         ], MID),
        ("Original creative direction",
         "Flagship brand campaigns without human leadership = brand damage",
         [
             "Toys R Us Cannes 2024 — sentiment swing −10 pp",
             "Coca-Cola Christmas 2024 — a cool reception",
             "SI by-line — an investigative scandal",
         ], LIGHT),
        ("Coherent long-form narrative",
         "50-min album · multi-act script — human for now",
         [
             "Character-continuity drift after 5-10 scenes",
             "Plot coherence across 90+ minutes — no",
             "Voice coherence across a whole album — no",
         ], GOLD),
    ]
    for i, (ztitle, zsubtitle, zpoints, accent) in enumerate(zones):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        filled_rect(s, x, card_y, card_w, 0.6, accent, radius=True, radius_adj=0.10)
        accent_text = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=card_y + 0.08, w=card_w - 0.30, h=0.45,
                 text=ztitle, size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 0.80, w=card_w - 0.36, h=0.8,
                 text=zsubtitle, size=12, italic=True, color=MID, line_spacing=1.30)
        for j, point in enumerate(zpoints):
            y = card_y + 1.80 + j * 0.85
            filled_rect(s, x + 0.18, y, 0.10, 0.10, accent, radius=True,
                        radius_adj=0.5)
            text_box(s, x=x + 0.40, y=y - 0.05, w=card_w - 0.55, h=0.75,
                     text=point, size=11, color=DEEP, line_spacing=1.35)
    # Footer
    text_box(s, x=0.55, y=6.85, w=12.3, h=0.4,
             text="Here — AI only as a supporting tool under human supervision; never as a replacement layer.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s34"))


def build_s35(p):
    """YouTube AI-превью."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="AI thumbnails on YouTube · 47.3% of creators opted out",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Empirical rejection by end users · Social Blade Creator Survey · Dec 2025",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 stat blocks
    stats = [
        ("47.3%", "of creators opted out of\nAI thumbnails", GOLD),
        ("−22% / −19%", "CTR drop:\ncreepy skin / text failure on mobile", LIGHT),
        ("−61.8%", "drop-off in the first 15 sec\npromise-content mismatch", MID),
    ]
    card_y = 2.0
    card_h = 3.0
    card_w = 3.95
    gap = 0.18
    start_x = 0.55
    for i, (big, label, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        filled_rect(s, x, card_y, card_w, 0.6, color, radius=True, radius_adj=0.10)
        accent_text = WHITE if color != GOLD else DEEP
        text_box(s, x=x + 0.15, y=card_y + 0.08, w=card_w - 0.30, h=0.45,
                 text=f"#{i+1}", size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.15, y=card_y + 0.85, w=card_w - 0.30, h=1.0,
                 text=big, size=36, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=x + 0.20, y=card_y + 1.95, w=card_w - 0.40, h=1.0,
                 text=label, size=12, color=DEEP, italic=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.40)
    # Lesson (expanded to full width, fills freed space below 3 stat cards)
    lesson_box(s, 0.55, 5.30, 12.3, 1.85,
               "End users notice and punish AI corner-cutting. Brand-trust risk is not theoretical — measurable CTR + drop-off + retention. If a product relies on AI visuals for end users, measure not only the cost of generation — measure CTR, retention, brand attitude.")
    speaker_notes(s, load_notes("s35"))


def build_s36(p):
    build_section_divider(p, here_idx=5, title="What the engineer should do",
                          frame_phrase="Action checklist · 5-question · the lecture's main takeaway.",
                          notes_slide_id="s36")


def build_s37(p):
    """5-question checklist."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="5-question checklist before AI in a creative project",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Apply it BEFORE the start, not after. If even one \"no/risk\" — reconsider the approach.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 5 questions vertically
    questions = [
        ("1.", "Is the tool's training data licensed?",
         "Firefly = yes (Adobe Stock + licensed content) · SD/Midjourney = Andersen-class risks"),
        ("2.", "Output-similarity check against protected content?",
         "NYT verbatim-quotation risk · a technical control is mandatory"),
        ("3.", "Consent for voice/likeness — if applicable?",
         "ScarJo · SAG-AFTRA Digital Replicas · Korea — always consent for real people"),
        ("4.", "IP-clean tools for commercial use?",
         "End-to-end pipeline: licensed corpus + similarity check + disclosure"),
        ("5.", "Brand-trust risk for flagship/historical campaigns?",
         "Coca-Cola · Toys R Us · SI — measurable spend of brand capital"),
    ]
    q_x = 0.55
    q_w = 8.5
    q_y = 2.0
    q_h = 0.95
    q_gap = 0.10
    for i, (num, qtitle, qbody) in enumerate(questions):
        y = q_y + i * (q_h + q_gap)
        ocean_box(s, q_x, y, q_w, q_h)
        # Number badge
        filled_rect(s, q_x + 0.10, y + 0.12, 0.7, q_h - 0.24, MID,
                    radius=True, radius_adj=0.20)
        text_box(s, x=q_x + 0.10, y=y + 0.12, w=0.7, h=q_h - 0.24,
                 text=num, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=q_x + 0.95, y=y + 0.10, w=q_w - 1.05, h=0.40,
                 text=qtitle, size=14, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=q_x + 0.95, y=y + 0.50, w=q_w - 1.05, h=0.40,
                 text=qbody, size=11, italic=True, color=SLATE, line_spacing=1.20)
    # Right: decision branch
    ocean_box(s, 9.25, 2.0, 3.7, 5.25, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.8)
    text_box(s, x=9.45, y=2.15, w=3.4, h=0.4,
             text="IF \"NO/RISK\" — 3 OPTIONS",
             size=11, bold=True, color=GOLD)
    options = [
        ("A.", "Don't use AI", "A non-AI alternative", MID),
        ("B.", "Mitigate structurally",
         "Licensing + similarity check + consent infrastructure + brand accounting", LIGHT),
        ("C.", "Accept the risk explicitly",
         "A documented business decision + calibrated mitigation. NOT implicitly.", TEAL),
    ]
    for i, (oletter, ot, ob, oc) in enumerate(options):
        y = 2.65 + i * 1.50
        filled_rect(s, 9.45, y, 0.5, 0.5, oc, radius=True, radius_adj=0.20)
        text_box(s, x=9.45, y=y, w=0.5, h=0.5, text=oletter,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=10.05, y=y, w=2.85, h=0.45,
                 text=ot, size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=9.45, y=y + 0.55, w=3.4, h=0.85,
                 text=ob, size=11, italic=True, color=SLATE, line_spacing=1.30)
    speaker_notes(s, load_notes("s37"))


def build_s38(p):
    """Q&A."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Left: big Q&A
    text_box(s, x=0.55, y=1.5, w=5.5, h=4.5, text="Q&A?",
             size=200, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # Right: central question recap + резервные вопросы
    ocean_box(s, 6.30, 1.5, 6.65, 5.0)
    text_box(s, x=6.50, y=1.65, w=6.3, h=0.4,
             text="CENTRAL QUESTION — A REMINDER",
             size=11, bold=True, color=TEAL)
    text_box(s, x=6.50, y=2.10, w=6.3, h=1.4,
             text="What has AI done to the creative industry by 2026 — and where to say \"no\"?",
             size=20, italic=True, bold=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=6.50, y=3.65, w=6.3, h=0.4,
             text="BACKUP QUESTIONS",
             size=11, bold=True, color=TEAL)
    prompts = [
        "Where are the limits of \"fair use\" in AI training?",
        "Sora vs Lionsgate — where is Hollywood in 2030?",
        "The Mintsifry bill — what changes for RU engineers?",
    ]
    for i, prompt in enumerate(prompts):
        y = 4.10 + i * 0.65
        filled_rect(s, 6.50, y, 6.3, 0.55, SURFACE,
                    stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.20)
        text_box(s, x=6.65, y=y, w=6.1, h=0.55,
                 text=f"·  {prompt}", size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s38"))


def build_s39(p):
    """Closing с hero-image (X-62 VISTA) — bridge к Лекции 9 (AI в авиакосмосе и обороне)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # ---- LEFT: thank you (compacted) ----
    ocean_box(s, 0.45, 1.30, 5.40, 4.95)
    text_box(s, x=0.65, y=1.50, w=5.05, h=0.45,
             text="THANK YOU", size=18, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT)
    text_box(s, x=0.65, y=2.15, w=5.05, h=1.50,
             text="for your attention",
             size=40, bold=True, color=DEEP, line_spacing=1.10)
    text_box(s, x=0.65, y=3.95, w=5.05, h=0.45,
             text="Checklist and sources — via the QR below",
             size=13, italic=True, color=MID, line_spacing=1.20)
    # QR placeholder
    filled_rect(s, 0.65, 4.75, 1.0, 1.0, WHITE, stroke=DEEP, stroke_pt=1.5)
    text_box(s, x=0.65, y=4.75, w=1.0, h=1.0, text="QR",
             size=18, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=1.80, y=4.75, w=3.90, h=1.0,
             text="Full bibliography of the lecture's sources (chapter.md + research dossier)",
             size=10, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.30)
    # ---- RIGHT: HERO BRIDGE IMAGE (X-62 VISTA — Лекция 9 foreshadow) ----
    # Image area: 7.0×3.95 (preserving 16:9 aspect of X-62 photo) ≈ 27.7 sq in ≈ 28%
    # Combined with right text panel — image+caption block ≈ 40% площади
    ocean_box(s, 6.10, 0.55, 6.75, 4.30, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, str(ASSETS / "screenshots/s39-x62-vista.jpg"),
              x=6.20, y=0.65, w=6.55, h=4.10)
    # Attribution chip under image
    text_box(s, x=6.20, y=4.90, w=6.55, h=0.30,
             text="X-62 VISTA · USAF Test Pilot School · DARPA ACE AI dogfight 2023 · Wikimedia Commons",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # Next lecture text block under hero
    text_box(s, x=6.20, y=5.30, w=6.55, h=0.35,
             text="NEXT LECTURE", size=12, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT)
    text_box(s, x=6.20, y=5.65, w=6.55, h=0.75,
             text="AI in the aerospace industry and the defense sector",
             size=20, bold=True, color=DEEP, line_spacing=1.15, align=PP_ALIGN.LEFT)
    filled_rect(s, 6.20, 6.50, 0.05, 0.45, fill=GOLD)
    text_box(s, x=6.35, y=6.45, w=6.50, h=0.55,
             text="From a public attack surface — to the safety-critical",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT, line_spacing=1.20)
    speaker_notes(s, load_notes("s39"))


# ============================================================
# Main
# ============================================================

def main():
    p = setup_pres()
    # Сборка в порядке deck.yaml
    builders = [
        build_s01, build_s02, build_s03, build_s04, build_s05, build_s05a,
        build_s06, build_s07, build_s08, build_s09, build_s10, build_s10a,
        build_s11,
        build_s13, build_s14, build_s15, build_s16, build_s17,
        build_s19, build_s20, build_s21, build_s22, build_s23, build_s24,
        build_s25, build_s26, build_s27, build_s28, build_s29, build_s30,
        build_s31,
        build_s32, build_s33, build_s34, build_s35,
        build_s36, build_s37,
        build_s38, build_s39,
    ]
    for builder in builders:
        builder(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved {OUT} with {len(builders)} slides.")


if __name__ == "__main__":
    main()
