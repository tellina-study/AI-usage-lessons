"""
Full 41-slide build of Лекции 13 «AI в логистике и транспорте».

Source-of-truth: deck.yaml v1 + chapter v2 multi-part (~31k слов) + slides/*.md.

Issue #135 · downstream от chapter (book-first).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml ≈ 75 мин.

Lec-N-1 паттерн compliance: match lec-11 (cover + lecture-map + 5 section dividers +
выделенный Q&A; top progress bar только на dividers + cover).

Build via: python3 build_lec13.py — generates lec-13.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
RED_WARN = RGBColor(0xC0, 0x39, 0x2B)
ROADMAP = RGBColor(0xD9, 0xE2, 0xEC)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent
ASSETS = ROOT.parent / "assets"
OUT = ROOT / "lec-13.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *, size=16, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except: pass
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline_box(slide, x, y, w, h, lines, *, size=14, bold=False, color=DEEP,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
                  line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except: pass
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image_aspect(slide, path, x, y, w, h):
    """Add picture preserving aspect ratio (centered in box)."""
    p = Path(path)
    if not p.exists():
        # placeholder grey box
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        text_box(slide, x, y+h/2-0.3, w, 0.6, f"[отсутствует: {p.name}]",
                 size=11, color=SLATE, align=PP_ALIGN.CENTER)
        return None
    try:
        with Image.open(p) as img:
            iw, ih = img.size
        img_ratio = iw / ih
        box_ratio = w / h
        if img_ratio > box_ratio:
            # width-constдождьed — fit to box width
            new_w = w
            new_h = w / img_ratio
            cx = x
            cy = y + (h - new_h) / 2
        else:
            # height-constдождьed — fit to box height
            new_h = h
            new_w = h * img_ratio
            cx = x + (w - new_w) / 2
            cy = y
        return slide.shapes.add_picture(str(p), Inches(cx), Inches(cy),
                                         width=Inches(new_w), height=Inches(new_h))
    except Exception as e:
        print(f"Image fail {path}: {e}")
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        return None


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


SECTIONS = ["Контролируемое", "Магистраль", "Город+миля", "Чрёзв.ситуация", "Замыкание"]


def roadmap_bar(slide, current_section):
    """5-section дорожная карта bar at top of section dividers + cover."""
    bar_y = 0.4
    bar_h = 0.32
    total_w = 12.33
    seg_w = total_w / 5
    for i, name in enumerate(SECTIONS):
        x = 0.5 + i * seg_w
        is_active = (i == current_section)
        rectangle(slide, x, bar_y, seg_w - 0.05, bar_h,
                  fill=GOLD if is_active else ROADMAP)
        text_box(slide, x, bar_y + 0.04, seg_w - 0.05, bar_h - 0.04,
                 f"{i+1}. {name}", size=10, bold=is_active,
                 color=DEEP if is_active else SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, text, *, y=7.05):
    text_box(slide, 0.5, y, 12.33, 0.35, text,
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)


def attribution(slide, text, *, x=0.5, y=6.95, w=12.33):
    text_box(slide, x, y, w, 0.3, text,
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)


def section_divider(p, num, title, items, section_idx):
    """Generic section divider with large number + title + list of items."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=section_idx)
    # Decorative large number
    text_box(slide, 0.5, 1.3, 5.0, 5.0, str(num),
             size=320, bold=True, color=ROADMAP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    # Title
    multiline_box(slide, 5.7, 2.0, 7.1, 1.3, [
        ("Раздел", {"size": 18, "color": LIGHT, "bold": True}),
        (title, {"size": 38, "bold": True, "color": DEEP}),
    ], line_spacing=1.1)
    # Item list with bullets
    multiline_box(slide, 5.7, 3.6, 7.1, 3.0, [(f"·  {x}", {"size": 14, "color": MID}) for x in items],
                  line_spacing=1.5)
    return slide


def assertion_visual(p, assertion, body_lines, *, image_path=None, attribution_text=None,
                     gold_callout=None, footer_text=None):
    """Generic assertion_visual layout: title top, body+image middle, footer."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    # Assertion title (max 2 lines)
    text_box(slide, 0.5, 0.4, 12.33, 1.3, assertion,
             size=24, bold=True, color=DEEP, line_spacing=1.15)
    return slide


# ========== SECTION 0 ==========

def s01_hero_three_pictures(p):
    """s01 — hero hook: Waymo Jaguar I-Pace SF + three time-cards.
    Hero ≥40% area: image 4:3 ratio, box 7.5×5.7 → fit 7.5×5.6=42.0 sq in = 42.0% of canvas."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    # Hero image left
    img_path = ASSETS / "screenshots" / "s01-waymo-jaguar-sf-dllu.jpg"
    add_image_aspect(slide, img_path, 0.5, 0.4, 7.5, 5.7)
    # Caption under hero
    attribution(slide, "Waymo Jaguar I-Pace в Сан-Франциско · 2023 · Wikimedia Commons · CC-BY-SA",
                x=0.5, y=6.15, w=7.5)
    # Title right
    multiline_box(slide, 8.3, 0.5, 4.7, 5.4, [
        ("Три картинки", {"size": 22, "bold": True, "color": DEEP}),
        ("рядом —", {"size": 22, "bold": True, "color": DEEP}),
        ("и один вопрос", {"size": 22, "bold": True, "color": MID}),
        ("", {"size": 8}),
        ("Декабрь 2024 · Cruise:", {"size": 12, "bold": True, "color": LIGHT}),
        ("10 миллиардов → 0", {"size": 14, "bold": True, "color": RED_WARN}),
        ("", {"size": 4}),
        ("Март 2026 · Waymo:", {"size": 12, "bold": True, "color": LIGHT}),
        ("500 000 поездок/неделю", {"size": 14, "bold": True, "color": GOLD}),
        ("", {"size": 4}),
        ("Февраль 2026 · Tesla:", {"size": 12, "bold": True, "color": LIGHT}),
        ("14 ДТП за 8 мес.", {"size": 14, "bold": True, "color": DEEP}),
    ], line_spacing=1.1)
    # Central question gold callout
    rounded_box(slide, 0.5, 6.0, 12.33, 1.0, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 0.8, 6.15, 11.8, 0.8,
             "Дело не в технологии — дело в среде, в которой система разворачивалась, и в дисциплине ODD.",
             size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "См. slides/s01-hook-three-pictures.md speaker notes.")


def s02_cover(p):
    """s02 — cover with decorative «13» and meta."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=-1)  # всех grey (no active)
    # Large decorative «13»
    text_box(slide, 0.5, 1.2, 4.0, 4.5, "13",
             size=240, bold=True, color=ROADMAP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    # Title
    multiline_box(slide, 4.5, 1.5, 8.3, 2.5, [
        ("Лекция 13", {"size": 22, "bold": True, "color": LIGHT}),
        ("AI в логистике", {"size": 40, "bold": True, "color": DEEP}),
        ("и транспорте", {"size": 40, "bold": True, "color": DEEP}),
    ], line_spacing=1.05)
    # Meta
    multiline_box(slide, 4.5, 4.4, 8.3, 1.0, [
        ("Модуль 3", {"size": 16, "color": LIGHT}),
        ("Студенты-инженеры 3 курса", {"size": 16, "color": LIGHT}),
    ])
    # Обучение outcomes (descriptive prose, нет LO codes)
    rounded_box(slide, 4.5, 5.5, 8.3, 1.6)
    multiline_box(slide, 4.7, 5.65, 8.0, 1.4, [
        ("Вы научитесь:", {"size": 14, "bold": True, "color": MID}),
        ("· Различать пять уровней автономии в логистике и транспорте.", {"size": 12, "color": DEEP}),
        ("· Критически оценивать вендорские заявления об AV/AI-системах.", {"size": 12, "color": DEEP}),
        ("· Применять регуляторные критерии «когда AI не работает» + альтернативы.", {"size": 12, "color": DEEP}),
    ], line_spacing=1.2)
    footer(slide, "Курс «Применение AI в инженерии» · 2026")
    add_notes(slide, "См. slides/s02-cover.md speaker notes.")


def s03_lecture_map(p):
    """s03 — 5 horizontal section cards."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Маршрут лекции — пять разделов",
             size=28, bold=True, color=DEEP)
    sections = [
        ("1", "Склад / порт", "6 слайдов", "Склад · порт · рельсы · границы уровня 1", LIGHT),
        ("2", "Магистраль", "10 слайдов", "Aurora · Mobileye · КамАЗ · UPS ORION · банкротства AV-грузоперевозки", MID),
        ("3", "Город + миля", "10 слайдов", "Waymo · Apollo Go · Pony.ai · Tesla · Cruise · Uber · Tesla NHTSA · Starship · Zipline", TEAL),
        ("4", "Чёрный лебедь", "7 слайдов", "Хуситы · Suez · COVID · дефицит дальнобойщиков · 5 критериев · альтернативы", GOLD),
        ("5", "Замыкание", "2 слайда", "Q&A · мост к Лекции 14", DEEP),
    ]
    card_w = 2.42
    разрыв = 0.05
    x0 = 0.5
    y = 1.4
    for i, (num, title, dur, desc, accent) in enumerate(sections):
        x = x0 + i * (card_w + разрыв)
        rounded_box(slide, x, y, card_w, 5.2)
        # Number circle on top-left
        circle(slide, x + 0.15, y + 0.2, 0.55, 0.55, fill=accent)
        text_box(slide, x + 0.15, y + 0.2, 0.55, 0.55, num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title (on its own line, full card width)
        text_box(slide, x + 0.15, y + 0.85, card_w - 0.3, 0.5, title,
                 size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        # Slides count
        text_box(slide, x + 0.15, y + 1.4, card_w - 0.3, 0.35, dur,
                 size=11, italic=True, color=LIGHT)
        # Description
        text_box(slide, x + 0.15, y + 1.85, card_w - 0.3, 3.2, desc,
                 size=12, color=DEEP, line_spacing=1.3)
    footer(slide, "Спуск по лестнице среды: от контролируемая к чёрному лебедю")
    add_notes(slide, "См. slides/s03-lecture-map.md speaker notes.")


def s04_glossary(p):
    """s04 — glossary mini: 10 terms × 2 columns."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Десять аббревиатур, без которых дальше не пройти",
             size=26, bold=True, color=DEEP)
    # Sub-title
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Глоссарий обязательных терминов — расшифровка + роль в лекции",
             size=14, italic=True, color=LIGHT)
    # Left column
    left = [
        ("SAE J3016", "Общество автомобильных инженеров. Стандарт уровней автомобильной автономии L0-L5."),
        ("ODD", "Операционный Design Domain — область штатной эксплуатации. Главная инженерная дисциплина."),
        ("AV", "Автономный Автомобиль — автономный автомобиль. Зонтичный термин."),
        ("AMR", "Автономный мобильный робот — автономный мобильный робот, чаще всего складской."),
        ("HD-map", "карта высокой точности — детальная карта с точностью ~10 см."),
    ]
    right = [
        ("OR", "Исследование операций — операционные исследования. НЕ AI. Канонические инструменты — Gurobi, CPLEX, OR-Tools."),
        ("TSP / VRP", "задача коммивояжёра / Автомобиль задача маршрутизации. UPS ORION экономит ~$300M/год через OR."),
        ("EOQ", "Economic размер заказа (Харрис 1913). Формула трёх параметров; для большой доли SKU не хуже ML."),
        ("SGO", "NHTSA постановление NHTSA on отчётность об авариях. База Tesla/Waymo/Cruise смертельный случай цифры."),
        ("HITL/HOOL/HOTL", "Human In/On/вне-цикла. Три уровня участия человека (Лекция 9)."),
    ]
    y0 = 1.7
    row_h = 1.0
    for i, (term, defn) in enumerate(left):
        rounded_box(slide, 0.5, y0 + i * row_h, 6.1, row_h - 0.05)
        text_box(slide, 0.7, y0 + i * row_h + 0.1, 1.8, 0.4, term,
                 size=14, bold=True, color=MID)
        text_box(slide, 0.7, y0 + i * row_h + 0.45, 5.6, row_h - 0.5, defn,
                 size=11, color=DEEP, line_spacing=1.2)
    for i, (term, defn) in enumerate(right):
        rounded_box(slide, 6.75, y0 + i * row_h, 6.1, row_h - 0.05)
        text_box(slide, 6.95, y0 + i * row_h + 0.1, 1.8, 0.4, term,
                 size=14, bold=True, color=TEAL)
        text_box(slide, 6.95, y0 + i * row_h + 0.45, 5.6, row_h - 0.5, defn,
                 size=11, color=DEEP, line_spacing=1.2)
    footer(slide, "Глоссарий — если хоть три термина незнакомы, зафиксируйте сейчас")
    add_notes(slide, "См. slides/s04-glossary.md speaker notes.")


def s05_keystone_ladder(p):
    """s05 — keystone: 5-step environment ladder."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Лестница среды: главный предиктор успеха AI в логистике",
             size=24, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.1, 12.33, 0.4,
             "Пять уровней — от контролируемого склада к чёрному лебедю. Каждый раздел = мотивированный шаг по этой оси.",
             size=13, italic=True, color=LIGHT)
    # 5 horizontal steps
    levels = [
        ("1", "Контролируемое", "Symbotic+Walmart 400 APD\nAmazon 1M+ роботов\nKONUX железная дорога PdM", "Капитальная интенсивность $$$", LIGHT),
        ("2", "Магистраль", "Aurora Dallas-Houston\nMobileye L3 eyes-off\nКамАЗ М-11 «Нева»", "$20B+ сожжено невыжившие", MID),
        ("3", "Городской robotaxi", "Waymo 500K/неделю\nApollo Go 22 города\nPony.ai прибыль на машину", "Cruise: $10B → 0", TEAL),
        ("4", "Последняя миля", "Starship 9M+ кампусы\nZipline 100M миль Africa\nNuro разворот 2024", "Дроны в городе США заблокированы", GOLD),
        ("5", "Чёрный лебедь", "Хуситы 2024 (-90%)\nSuez 2021 (12% торговли)\nCOVID 2020 обвал", "AI не работает — нужен человек", RED_WARN),
    ]
    step_w = 2.42
    gap = 0.05
    x0 = 0.5
    y = 1.7
    for i, (num, name, examples, fail, accent) in enumerate(levels):
        x = x0 + i * (step_w + gap)
        h = 4.5
        # Box height grows with уровень для visual progression
        rounded_box(slide, x, y, step_w, h, stroke=accent, stroke_w=2.0)
        # Number circle
        circle(slide, x + step_w/2 - 0.4, y + 0.15, 0.8, 0.8, fill=accent)
        text_box(slide, x + step_w/2 - 0.4, y + 0.15, 0.8, 0.8, num,
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        text_box(slide, x + 0.1, y + 1.15, step_w - 0.2, 0.4, name,
                 size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        # Examples
        text_box(slide, x + 0.15, y + 1.65, step_w - 0.3, 2.0, examples,
                 size=10, color=DEEP, line_spacing=1.35)
        # Провал metric (bottom, gold or red)
        rectangle(slide, x + 0.1, y + h - 0.95, step_w - 0.2, 0.8, fill=GOLD_TINT)
        text_box(slide, x + 0.15, y + h - 0.85, step_w - 0.3, 0.65, fail,
                 size=10, italic=True, bold=True, color=DEEP, line_spacing=1.25)
    # Bottom takeaway
    rounded_box(slide, 0.5, 6.4, 12.33, 0.7, fill=GOLD_TINT, stroke=GOLD)
    text_box(slide, 0.7, 6.45, 11.9, 0.6,
             "Среда определяет AI. Выжившие уважают среду + остаются в узком ODD + не переобещают.",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "См. slides/s05-keystone-environment-ladder.md speaker notes.")


# ========== SECTION 1 ==========

def s06_section1(p):
    """s06 — divider Раздел 1."""
    section_divider(p, 1, "Контролируемое",
                    ["Склад: Symbotic+Walmart, Amazon Robotics, AMR (Locus / GreyOrange / Geek+)",
                     "Порт: ABB, Konecranes, ZPMC — и ILA-страйки 2024",
                     "Железная дорога: KONUX + Deutsche Bahn",
                     "Границы уровня 1: гумано-робот ажиотажа, безлюдный миф, капитал интенсивность, распределение сдвиг"],
                    section_idx=0)


def s07_symbotic_walmart(p):
    """s07 — Symbotic + Walmart кейс."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Symbotic + Walmart: 400 APD-центров,\nболее 5 миллиардов портфель заказов",
             size=24, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image left
    img_path = ASSETS / "screenshots" / "s07-symbotic-logo.jpg"
    rounded_box(slide, 0.5, 1.9, 5.5, 4.3)
    add_image_aspect(slide, img_path, 0.7, 2.1, 5.1, 3.9)
    attribution(slide, "Symbotic · Wikimedia Commons", x=0.5, y=6.25, w=5.5)
    # Right panel facts
    rounded_box(slide, 6.3, 1.9, 6.5, 4.3)
    multiline_box(slide, 6.5, 2.05, 6.1, 4.0, [
        ("Цифры", {"size": 14, "bold": True, "color": MID}),
        ("400 APD-центров — Walmart обязался развернуть Symbotic (январь 2025)", {"size": 12, "color": DEEP}),
        ("Более 5 миллиардов долларов портфель заказов — экстраординарный сигнал", {"size": 12, "color": DEEP}),
        ("(4× годовой выручки вперёд)", {"size": 11, "italic": True, "color": LIGHT}),
        ("$618 миллионов — выручка Symbotic Q4 FY2025", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Что Symbotic делает технически", {"size": 14, "bold": True, "color": MID}),
        ("Паллетизация · депаллетизация · отбор · сортировка", {"size": 12, "color": DEEP}),
        ("Комбинация классических роботов + мобильных платформ + CV", {"size": 12, "color": DEEP}),
        ("ML для распознавания SKU, основное управление — детерминированные правила + WMS", {"size": 11, "italic": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Режим провала", {"size": 14, "bold": True, "color": GOLD}),
        ("Капитальная интенсивность — десятки миллионов на APD-центр", {"size": 12, "color": DEEP}),
        ("Малые ритейлеры не могут позволить · OEM привязка к вендору риск", {"size": 12, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Уровень 1 лестницы среды · самая контролируемая среда логистики")
    add_notes(slide, "См. slides/s07-symbotic-walmart.md speaker notes.")


def s08_amazon_robotics(p):
    """s08 — Amazon Robotics: 4 поколения."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Amazon: более 1 миллиона роботов в фулфилменте.\nSparrow → Vulcan: путь от только-зрение к тактильному.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image left
    img_path = ASSETS / "screenshots" / "s08-amazon-warehouse-robot-2020.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 5.5, 4.3)
    attribution(slide, "Amazon склад робот · 2020 · Wikimedia · CC-BY-SA", x=0.5, y=6.25, w=5.5)
    # 4 робот panels right (2×2)
    роботs = [
        ("Sparrow", "2022", "Захват из ячейки манипулятор. Зрение + глубина. Ограничения: текстиль, прозрачная упаковка.", LIGHT),
        ("Sequoia", "2023", "Контейнерное хранение. Мехатроника + WMS-логика, не «AI» в строгом смысле.", MID),
        ("Proteus", "2022", "Первый полностью автономный AMR в фулфилменте. SLAM в условиях движущихся людей.", TEAL),
        ("Vulcan", "2024-25", "+ тактильная чувствительность. Ответ на ограничения Sparrow (текстиль, зеркала).", GOLD),
    ]
    panel_w = 3.1
    panel_h = 2.1
    for i, (name, year, desc, color) in enumerate(роботs):
        col = i % 2
        row = i // 2
        x = 6.3 + col * (panel_w + 0.1)
        y = 1.9 + row * (panel_h + 0.1)
        rounded_box(slide, x, y, panel_w, panel_h, stroke=color, stroke_w=2)
        text_box(slide, x + 0.15, y + 0.1, panel_w - 0.3, 0.4, name,
                 size=15, bold=True, color=color)
        text_box(slide, x + 0.15, y + 0.5, panel_w - 0.3, 0.3, year,
                 size=10, italic=True, color=SLATE)
        text_box(slide, x + 0.15, y + 0.85, panel_w - 0.3, panel_h - 1.0, desc,
                 size=11, color=DEEP, line_spacing=1.25)
    # Bottom callout: lesson
    rounded_box(slide, 6.3, 6.4, 6.5, 0.7, fill=GOLD_TINT, stroke=GOLD)
    text_box(slide, 6.5, 6.45, 6.2, 0.6,
             "Зрение-стек не универсален. Multi-modal сенсор — урок 3 лет.",
             size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    footer(slide, "Amazon ~1M роботов 2025 · крупнейший складской автоматизатор в мире")
    add_notes(slide, "См. slides/s08-amazon-robotics.md speaker notes.")


def s09_amr_locus(p):
    """s09 — AMR стек Locus + GreyOrange + Geek+ + worker отпор."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "AMR-стек: Locus 5 миллиардов операций отбора.\nНо AMR — не из коробки (требует развёртывания), пушбэк рабочих по нагрузке — реален.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image left
    img_path = ASSETS / "screenshots" / "s09-item-picking-robot.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 5.5, 4.3)
    attribution(slide, "Item Picking Робот · Wikimedia · CC-BY-SA", x=0.5, y=6.25, w=5.5)
    # 3 player cards right
    players = [
        ("Locus Robotics", "Совместный отбор. 5+ миллиардов операций отбора, 350+ развёртываний.\nFedEx · GXO · DHL.", LIGHT),
        ("GreyOrange Butler", "Товары-to-person. Рак-полка движется к человеку. Плотность хранения растёт ×1,5-2.", MID),
        ("Geek+ (Китай)", "Лидер по объёмам в Азии. Hybrid совместный + «товар-к-человеку». Расширение в Европу.", TEAL),
    ]
    py = 1.9
    for name, desc, color in players:
        rounded_box(slide, 6.3, py, 6.5, 1.4, stroke=color, stroke_w=2)
        text_box(slide, 6.5, py + 0.1, 6.1, 0.4, name,
                 size=14, bold=True, color=color)
        text_box(slide, 6.5, py + 0.5, 6.1, 0.85, desc,
                 size=11, color=DEEP, line_spacing=1.25)
        py += 1.5
    # Bottom worker отпор callout
    rounded_box(slide, 6.3, 6.4, 6.5, 0.7, fill=GOLD_TINT, stroke=GOLD)
    text_box(slide, 6.5, 6.45, 6.2, 0.6,
             "Worker нагрузка отпор: 50-100 pick/час → 200-250. Toyota Jidoka — augment, не заменить.",
             size=11, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    footer(slide, "AMR не готовый к работе · развёртывание 3-9 месяцев · перепланировка склада")
    add_notes(slide, "См. slides/s09-amr-locus-greyorange.md speaker notes.")


def s10_port_automation(p):
    """s10 — Port automation + ILA strikes."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Maasvlakte II · Long Beach LBCT · Yangshan — автоматизированы.\nILA-страйки 2024: технология упирается в трудовую политику.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero: Long Beach terminal
    img_path = ASSETS / "screenshots" / "s10b-long-beach-terminal.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.0, 4.3)
    attribution(slide, "Terminal Island, Long Beach · Wikimedia · CC-BY-SA", x=0.5, y=6.25, w=6.0)
    # Right panel
    rounded_box(slide, 6.8, 1.9, 6.0, 4.3)
    multiline_box(slide, 7.0, 2.05, 5.6, 4.0, [
        ("Главные поставщики", {"size": 14, "bold": True, "color": MID}),
        ("ABB · STS краны semi-automated (Швейцария)", {"size": 11, "color": DEEP}),
        ("Konecranes · автоматизированный контейнер обработка (Финляндия)", {"size": 11, "color": DEEP}),
        ("ZPMC · ~70% мирового рынка STS-кранов (Шанхай)", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Что НЕ автоматизировано", {"size": 14, "bold": True, "color": MID}),
        ("Quay краны — финальное позиционирование = человек/remote", {"size": 11, "color": DEEP}),
        ("Lashing (закрепление контейнеров) — ручная работа", {"size": 11, "color": DEEP}),
        ("Inspection · таможня — диспетчеры-люди", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("ILA-страйки октябрь 2024", {"size": 14, "bold": True, "color": GOLD}),
        ("~85 000 работников · до $5 миллиардов/день ущерба", {"size": 11, "color": DEEP}),
        ("US East+Gulf Coast — $75-100 миллиардов грузопоток/год", {"size": 11, "color": DEEP}),
        ("Урок: технология сталкивается с трудовой политика", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.2)
    footer(slide, "Maasvlakte II Rotterdam · LBCT Long Beach · Yangshan Shanghai")
    add_notes(slide, "См. slides/s10-port-automation.md speaker notes.")


def s11_rail_konux(p):
    """s11 — Rail KONUX + границы уровня 1."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "KONUX + Deutsche Bahn: PdM стрелок зрело.\nГумано-робот и безлюдный склад — пока миф.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero left: ICE tдождь
    img_path = ASSETS / "screenshots" / "s11-deutsche-bahn-ice.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 5.5, 2.3)
    attribution(slide, "Deutsche Bahn ICE Testfahrt · Wikimedia · CC-BY-SA", x=0.5, y=4.25, w=5.5)
    # Below: железная дорога стрелка
    img_path = ASSETS / "screenshots" / "s11b-rail-switch.jpg"
    add_image_aspect(slide, img_path, 0.5, 4.6, 5.5, 2.0)
    attribution(slide, "Railway стрелка lever · Wikimedia · CC-BY-SA", x=0.5, y=6.65, w=5.5)
    # Right: 3 границы
    borders = [
        ("Гумано-робот ажиотажа", "Tesla Optimus / Figure 02 / UBTech — на 2026 в фазе исследования / пилот. Контрфактив: Amazon Sparrow $50-150K vs гуманоид $100-300K.", GOLD),
        ("Капитальная интенсивность", "Один APD-центр Symbotic — десятки миллионов. ZPMC-кран — $10M+. Maasvlakte II — €2 миллиарда. Доступно только крупным.", MID),
        ("Без-света склад миф", "Работает только в узких категориях (автозапчасти, банковский архив). Широкий SKU с миллионами артикулов — пока нет.", LIGHT),
    ]
    py = 1.9
    for title, desc, color in borders:
        rounded_box(slide, 6.3, py, 6.5, 1.5, stroke=color, stroke_w=2)
        text_box(slide, 6.5, py + 0.1, 6.1, 0.4, title,
                 size=14, bold=True, color=color)
        text_box(slide, 6.5, py + 0.5, 6.1, 0.95, desc,
                 size=11, color=DEEP, line_spacing=1.25)
        py += 1.6
    footer(slide, "Педагогически: тот же пилот чистилище как в Лекции 11 — adoption ограничен бюрократией")
    add_notes(slide, "См. slides/s11-rail-konux-limits.md speaker notes.")


def s12_discrete_failure_matrix(p):
    """s12 — провал matrix Раздел 1."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Четыре границы уровня 1: капитал интенсивность, OEM привязка к вендору,\nбез-света миф, сезонный распределение сдвиг.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Matrix 4 rows × 3 cols
    cols = ["Граница", "Симптом", "Что делать инженеру"]
    rows = [
        ("Капитальная интенсивность",
         "APD-центр Symbotic — десятки $M. ZPMC-кран — $10M+. Maasvlakte II — €2 миллиарда. Малые не могут.",
         "Ориентир для малого 3PL — конвейеры + штрих-код + ML-маршрутизация через SaaS (project44, FourKites)."),
        ("OEM привязка к вендору",
         "Walmart инвестирует миллиарды в Symbotic — стрелка на Knapp/Vanderlande/KION займёт годы и major re-инвестиции.",
         "При long-term контракт: выход strategy? proprietary vs стандартизированные компоненты? стоимость migration через 10 лет?"),
        ("Без-света склад миф",
         "Работает только в узких категориях. Широкий SKU с миллионами артикулов — долгий хвост краевые случаи (текстиль, returns).",
         "Спросить: «Какой % SKU без вмешательства человека?» Если 99% — попросить логи и список 1% краевые случаи."),
        ("Сезонный распределение сдвиг",
         "Чёрная пятница / Рождество / Halloween — каталог наполняется новыми категориями. Sparrow зрение-классификатор ошибается чаще.",
         "Спросить: частота ошибок на пиках vs обычное время? как часто переобучается? кто оплачивает retдождьing cost?"),
    ]
    y = 1.7
    # Header row
    col_w = [3.3, 4.5, 4.4]
    col_x = [0.5, 0.5+3.4, 0.5+3.4+4.6]
    for i, c in enumerate(cols):
        rectangle(slide, col_x[i], y, col_w[i], 0.5, fill=MID)
        text_box(slide, col_x[i]+0.1, y+0.05, col_w[i]-0.2, 0.4, c,
                 size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.55
    row_h = 1.2
    accent_colors = [LIGHT, MID, TEAL, GOLD]
    for i, (a, b, c) in enumerate(rows):
        ac = accent_colors[i]
        for j, val in enumerate([a, b, c]):
            rounded_box(slide, col_x[j], y, col_w[j], row_h,
                       stroke=ac if j == 0 else LIGHT, stroke_w=2 if j == 0 else 1)
            text_box(slide, col_x[j]+0.15, y+0.08, col_w[j]-0.3, row_h-0.15, val,
                     size=11 if j > 0 else 12, bold=(j == 0), color=DEEP, line_spacing=1.25)
        y += row_h + 0.05
    footer(slide, "Чек-лист «что спросить» при оценке любого решения уровня 1")
    add_notes(slide, "См. slides/s12-discrete-failure-matrix.md speaker notes.")


# ========== SECTION 2 ==========

def s13_section2(p):
    section_divider(p, 2, "Магистраль",
                    ["Aurora Инновации: первая беспилотный коммерческий US (май 2025)",
                     "Mobileye Шофёр L3 + КамАЗ-54901 М-11 «Нева»",
                     "UPS ORION: операционные исследования, не RL (канонический анти-хайп)",
                     "AV-хронология банкротств: Argo / Embark / TuSimple / Waymo Via / Starsky",
                     "Совокупно >$20 миллиардов сожжено на невыжившие",
                     "Выживший консолидация 10:1 — Waymo, Aurora, Mobileye"],
                    section_idx=1)


def s14_aurora(p):
    """s14 — Aurora first коммерческий L4."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Aurora Инновации: 1 мая 2025 — первая в США\nкоммерческая беспилотная грузоперевозка.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image
    img_path = ASSETS / "screenshots" / "s14-aurora-driverless-press.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.5, 4.3)
    attribution(slide, "Aurora водитель Class-8 грузовик · Aurora press, май 2025", x=0.5, y=6.25, w=6.5)
    # Right panel
    rounded_box(slide, 7.3, 1.9, 5.5, 4.3)
    multiline_box(slide, 7.5, 2.05, 5.1, 4.0, [
        ("Цифры", {"size": 14, "bold": True, "color": MID}),
        ("1 мая 2025 — старт беспилотный коммерческий", {"size": 11, "color": DEEP}),
        ("Маршрут Даллас-Хьюстон · ~240 миль I-45", {"size": 11, "color": DEEP}),
        ("~10 машин к концу 2025 · 100 план 2027", {"size": 11, "color": DEEP}),
        ("Платформа PACCAR Peterbilt + Aurora водитель", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Поэтапный ввод подход", {"size": 14, "bold": True, "color": MID}),
        ("Crawl 2018-2024: безопасность оператор, многолетние тесты", {"size": 11, "color": DEEP}),
        ("Walk 2023-2024: Dallas-Houston с оператор на отдельных рейсах", {"size": 11, "color": DEEP}),
        ("Run май 2025: беспилотный коммерческий единичный маршрут", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Цитата Chris Urmson (CEO)", {"size": 14, "bold": True, "color": GOLD}),
        ("«Новая эра грузоперевозок» — не «replacing X% drivers»", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Паттерн выживания: поэтапный ввод (ползком-шагом-бегом) · узкий ODD · не переобещают")
    add_notes(slide, "См. slides/s14-aurora-dallas-houston.md speaker notes.")


def s15_mobileye_kamaz(p):
    """s15 — Mobileye + КамАЗ."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Mobileye Шофёр L3 eyes-off + КамАЗ-54901 М-11:\nдве модели на полуструктурированной магистрали.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Left: Mobileye
    rounded_box(slide, 0.5, 1.9, 6.2, 4.5)
    text_box(slide, 0.7, 2.05, 5.8, 0.5, "Mobileye Шофёр",
             size=18, bold=True, color=MID)
    img_path = ASSETS / "screenshots" / "s15c-mobileye-zeekr-iaa.jpg"
    add_image_aspect(slide, img_path, 0.7, 2.6, 5.8, 2.3)
    attribution(slide, "Mobileye Zeekr · IAA Summit 2023, Munich · Wikimedia · CC-BY-SA",
                x=0.7, y=4.95, w=5.8)
    multiline_box(slide, 0.7, 5.3, 5.8, 1.0, [
        ("~300K SuperЗрение к 2025, цель 1,2M к 2026", {"size": 11, "color": DEEP}),
        ("Сначала-камера ADAS без дорогого LiDAR", {"size": 11, "color": DEEP}),
        ("L3 eyes-off на Polestar 4, Audi Q6, VW Touareg", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    # Right: КамАЗ
    rounded_box(slide, 6.9, 1.9, 6.0, 4.5)
    text_box(slide, 7.1, 2.05, 5.6, 0.4, "КамАЗ-54901 + Cognitive Pilot",
             size=18, bold=True, color=TEAL)
    text_box(slide, 7.1, 2.45, 5.6, 0.2,
             "Cognitive Pilot — российский разработчик стека восприятия для AV",
             size=9, italic=True, color=SLATE)
    img_path = ASSETS / "screenshots" / "s15b-kamaz-truck.jpg"
    add_image_aspect(slide, img_path, 7.1, 2.6, 5.6, 2.3)
    attribution(slide, "КамАЗ · Wikimedia Commons · CC-BY-SA", x=7.1, y=4.95, w=5.6)
    multiline_box(slide, 7.1, 5.3, 5.6, 1.0, [
        ("67 единиц на М-11 «Нева» (2024), 10 в коммерческий cargo", {"size": 11, "color": DEEP}),
        ("План 100 единиц на 2025 + М-12 + ЦКАД", {"size": 11, "color": DEEP}),
        ("В рамках ЭПР (экспериментального правового режима)", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    footer(slide, "Разные бизнес-модели · потребительская L3 vs коммерческий L4 на выделенный магистрали")
    add_notes(slide, "См. slides/s15-mobileye-kamaz.md speaker notes.")


def s16_ups_orion(p):
    """s16 — UPS ORION — канонический OR success."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "UPS ORION: 100M миль/год, $300-400M экономия/год.\nЭто операционные исследования, не глубокое обучение, не RL.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero left
    img_path = ASSETS / "screenshots" / "s16-ups-delivery-truck.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 5.5, 3.0)
    attribution(slide, "Typical UPS доставка грузовик · Wikimedia · CC-BY-SA",
                x=0.5, y=4.95, w=5.5)
    # Big gold callout цифры
    rounded_box(slide, 0.5, 5.3, 5.5, 1.6, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    multiline_box(slide, 0.7, 5.4, 5.1, 1.4, [
        ("$300-400 миллионов/год", {"size": 22, "bold": True, "color": GOLD}),
        ("savings для UPS через OR + эвристики", {"size": 12, "italic": True, "color": DEEP}),
        ("Парк ~125 000 машин · 100 миллионов миль/год экономии", {"size": 11, "color": DEEP}),
    ], line_spacing=1.2)
    # Right panel
    rounded_box(slide, 6.3, 1.9, 6.5, 5.0)
    multiline_box(slide, 6.5, 2.05, 6.1, 4.8, [
        ("Что под капотом", {"size": 14, "bold": True, "color": MID}),
        ("Исследование операций", {"size": 14, "bold": True, "color": TEAL}),
        ("ILP + эвристики + Автомобиль задача маршрутизации (VRP)", {"size": 11, "color": DEEP}),
        ("НЕ глубокое обучение · НЕ RL · НЕ GenAI", {"size": 12, "bold": True, "color": RED_WARN}),
        ("Используемые инструменты: Gurobi, CPLEX, OR-Tools", {"size": 11, "color": DEEP}),
        ("Алгоритмы датируются 1950-60-ми (теория графов, branch-and-bound)", {"size": 11, "italic": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Когда задавать поставщику вопросы", {"size": 14, "bold": True, "color": GOLD}),
        ("Какие сравнения с базовая линия OR (Gurobi / OR-Tools)?", {"size": 11, "color": DEEP}),
        ("Если нет — красный флаг.", {"size": 11, "italic": True, "color": DEEP}),
        ("Какой VRP-solver? «сквозной deep обучение» → объяснимость?", {"size": 11, "color": DEEP}),
        ("Stationary спрос или нет? Если да — EOQ может быть достаточным.", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Педагогически", {"size": 14, "bold": True, "color": LIGHT}),
        ("UPS ORION = канонический анти-хайп example. Простая задача — OR лучше ML.", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "FedEx, Maersk, ZIM, CMA CGM используют specialized OR-solvers для своих задач")
    add_notes(slide, "См. slides/s16-ups-orion-fundamental.md speaker notes.")


def s17_av_bankruptcy_timeline(p):
    """s17 — AV-грузоперевозки хронология банкротств 2020-2024."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "AV-грузоперевозки хронология банкротств 2020-2024:\nпять компаний на одной линии, серия катастроф.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Хронология horizontal line
    rectangle(slide, 1.0, 4.0, 11.3, 0.05, fill=LIGHT)
    # 5 событиеs on хронология
    событиеs = [
        ("Март 2020", "Starsky\nRobotics", "~$200M\nсожжено", "Первая волна\nсимуляция-к-реальности эссе"),
        ("Окт 2022", "Argo AI", "$7B\nсожжено", "Ford+VW тянуть финансирование\n$2,7B impairment"),
        ("Март 2023", "Embark\nTrucks", "$5,16B цель\ncap, 16 мес.", "От SPAC IPO\nдо банкротства"),
        ("2023", "Waymo Via", "Alphabet\nзакрыла", "Бесконечный капитал\nне нашёл прибыльный"),
        ("Янв 2024", "TuSimple", "91%+ shareholder\nvalue lost", "Nasdaq delisting\nUS-China напряжение"),
    ]
    n = len(событиеs)
    item_x_start = 1.0
    item_x_step = 11.3 / n
    for i, (date, name, money, desc) in enumerate(событиеs):
        x = item_x_start + i * item_x_step
        # Point on хронология
        circle(slide, x + item_x_step/2 - 0.15, 3.9, 0.3, 0.3, fill=GOLD if i in [1, 4] else MID)
        # Card above
        rounded_box(slide, x + 0.1, 1.7, item_x_step - 0.2, 2.1, stroke=GOLD if i in [1, 4] else LIGHT, stroke_w=2)
        text_box(slide, x + 0.15, 1.75, item_x_step - 0.3, 0.4, date,
                 size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
        text_box(slide, x + 0.15, 2.15, item_x_step - 0.3, 0.5, name,
                 size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(slide, x + 0.15, 2.7, item_x_step - 0.3, 0.5, money,
                 size=10, italic=True, color=RED_WARN, align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(slide, x + 0.15, 3.2, item_x_step - 0.3, 0.6, desc,
                 size=9, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.1)
    # Image bottom-left: Argo AI
    img_path = ASSETS / "screenshots" / "s17-argo-ai-vehicle-2021.jpg"
    add_image_aspect(slide, img_path, 0.5, 4.5, 5.5, 2.0)
    attribution(slide, "Argo AI (2021) · Wikimedia · CC-BY-SA", x=0.5, y=6.55, w=5.5)
    # Общие паттерны right
    rounded_box(slide, 6.3, 4.5, 6.5, 2.4)
    multiline_box(slide, 6.5, 4.6, 6.1, 2.2, [
        ("Общие паттерны невыжившие", {"size": 14, "bold": True, "color": MID}),
        ("• Капитал интенсивность без выручка ($1-7B сожжено до коммерции)", {"size": 11, "color": DEEP}),
        ("• SPAC IPO пузырь 2021-2022 (Embark, TuSimple, Aurora)", {"size": 11, "color": DEEP}),
        ("• симуляция-к-реальности разрыв (ML в demo ≠ ML на публичный дороги)", {"size": 11, "color": DEEP}),
        ("• Регуляторная неопределённость (NHTSA SGO, состояние мозаичный)", {"size": 11, "color": DEEP}),
        ("• Misaligned customer спрос (перевозчики хотели выделенный полосаs)", {"size": 11, "color": DEEP}),
        ("• Выживший консолидация 10:1 — выжили 3-4 из 30+", {"size": 11, "bold": True, "color": GOLD}),
    ], line_spacing=1.3)
    footer(slide, "См. также s22 Stefan Seltz-Axmacher разбор полётов")
    add_notes(slide, "См. slides/s17-av-bankruptcy-timeline.md speaker notes.")


def s18_cumulative_20b(p):
    """s18 — совокупно $20B+ chart."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Совокупно: >$20 миллиардов сожжено на невыжившие\nавтономный грузоперевозки + robotaxi 2017-2024.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Layout: [Company label | Bar | $ value] + desc row below
    companies = [
        ("Cruise (GM)", 10.0, "$10B операционный убыток · <$500M выручка · 2016-2024"),
        ("Argo AI", 7.0, "Ford $5B+ + VW $2,6B · 2017-2022"),
        ("TuSimple", 1.0, "IPO proceeds + private · 91% loss · 2017-2024"),
        ("Embark", 0.5, "SPAC цель $5,16B · 16 мес. жизни · 2020-2023"),
        ("Starsky", 0.2, "Первая волна · 2015-2020"),
        ("Waymo Via", 0.5, "Alphabet инвестиции · сегмент закрыт 2023"),
    ]
    max_amt = max(c[1] for c in companies)
    label_x = 0.5      # Company label column
    label_w = 1.8
    bar_x = 2.4        # Bar starts here
    bar_max_w = 8.0    # Bar drawing area
    value_x = 10.5     # $ value column
    value_w = 1.5
    row_h = 0.40       # Row total height for bar
    desc_h = 0.22      # Desc text row height
    row_gap = 0.03
    py = 1.85
    for name, amt, desc in companies:
        w = bar_max_w * amt / max_amt
        # Company label (left column, outside bar)
        text_box(slide, label_x, py, label_w, row_h,
                 name, size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        # Bar
        rectangle(slide, bar_x, py + 0.05, max(w, 0.15), row_h - 0.1,
                  fill=RED_WARN if amt >= 5 else LIGHT)
        # $ value (right column, outside bar)
        text_box(slide, value_x, py, value_w, row_h,
                 f"${amt:.1f}B" if amt >= 1 else f"${int(amt*1000)}M",
                 size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        # Desc row below
        text_box(slide, bar_x, py + row_h + row_gap, bar_max_w + 2.0, desc_h, desc,
                 size=10, italic=True, color=SLATE)
        py += row_h + desc_h + row_gap + 0.05
    # Total callout
    rounded_box(slide, 0.5, 6.15, 12.33, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    multiline_box(slide, 0.7, 6.20, 11.9, 0.75, [
        (">$20 миллиардов  ·  только невыжившие AV-грузоперевозки + robotaxi 2017-2024",
         {"size": 18, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("До $50 миллиардов с early-stage инвестицииs в Cruise, Zoox, Aurora ранние раунды",
         {"size": 12, "italic": True, "color": DEEP, "align": PP_ALIGN.CENTER}),
    ], align=PP_ALIGN.CENTER, line_spacing=1.2)
    footer(slide, "Lec-11 Pilot чистилище был $десятки тысяч/пилот · AV-провал в тысячи раз дороже")
    add_notes(slide, "См. slides/s18-cumulative-20b-burned.md speaker notes.")


def s19_survivor_consolidation(p):
    """s19 — выжившие vs выбывшие."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Выжившие vs выбывшие: 10:1 консолидация.\nWaymo, Aurora, Mobileye, Apollo Go — выжившие из 30+.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # 2-column matrix
    # Выжившие (left, gold)
    rounded_box(slide, 0.5, 1.9, 6.0, 4.6, stroke=GOLD, stroke_w=2)
    text_box(slide, 0.7, 2.0, 5.6, 0.5, "Выжившие (3-4 компании)",
             size=16, bold=True, color=GOLD)
    выжившие = [
        ("Waymo (Alphabet)", "Поэтапный ввод · HD-карта + LiDAR + удалённый опс · 500K rides/неделю"),
        ("Aurora Инновации", "One route at a время · Dallas-Houston коммерческий май 2025 · ~10 грузовиков"),
        ("Mobileye", "Сначала-камера ADAS · ~300K потребитель L2/L3 · spin-off Intel"),
        ("Apollo Go (Baidu)", "240M км глобально · 17M+ заказов · 22 города"),
    ]
    py = 2.55
    for name, desc in выжившие:
        text_box(slide, 0.7, py, 5.6, 0.35, name,
                 size=12, bold=True, color=DEEP)
        text_box(slide, 0.7, py + 0.35, 5.6, 0.55, desc,
                 size=10, color=DEEP, line_spacing=1.25)
        py += 1.0
    # Dropouts (right, red)
    rounded_box(slide, 6.7, 1.9, 6.1, 4.6, stroke=RED_WARN, stroke_w=2)
    text_box(slide, 6.9, 2.0, 5.7, 0.5, "Dropouts (>15 компаний)",
             size=16, bold=True, color=RED_WARN)
    выбывшие = [
        ("Argo AI", "$7B сгорело · Ford+VW отозвали финансирование одновременно"),
        ("Cruise (GM)", "$10B → 0 за 8 лет · October 2023 волочения инцидент"),
        ("TuSimple", "Delisting · китайский asset transfer · US-China напряжение"),
        ("Embark", "SPAC bust 16 месяцев от IPO"),
        ("Waymo Via", "Alphabet закрыла собственное грузоперевозки"),
        ("Starsky", "Первая волна 2020 · разрыв симуляция-к-реальности"),
    ]
    py = 2.55
    for name, desc in выбывшие:
        text_box(slide, 6.9, py, 5.7, 0.32, name,
                 size=11, bold=True, color=DEEP)
        text_box(slide, 6.9, py + 0.32, 5.7, 0.32, desc,
                 size=9, color=DEEP, line_spacing=1.2)
        py += 0.65
    # Паттерн выживания callout
    rounded_box(slide, 0.5, 6.55, 12.33, 0.55, fill=GOLD_TINT, stroke=GOLD)
    text_box(slide, 0.7, 6.6, 11.9, 0.45,
             "Паттерн выживания: поэтапный ввод (ползком-шагом-бегом) · узкий ODD · не переобещают · терпеливый капитал · уважение к среде.",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "См. slides/s19-survivor-consolidation.md speaker notes.")


def s20_trucker_shortage_false(p):
    """s20 — дефицит дальнобойщиков false framing."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "«AV решит дефицит водителей» — false framing.\nАрифметика не сходится на горизонте 2030.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Left: deficit number
    rounded_box(slide, 0.5, 1.9, 5.5, 4.5, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    text_box(slide, 0.7, 2.05, 5.1, 0.5, "ATA 2022 — пик",
             size=14, italic=True, color=LIGHT)
    text_box(slide, 0.7, 2.6, 5.1, 1.5, "78 000",
             size=80, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(slide, 0.7, 4.3, 5.1, 0.6, "дефицит водителей",
             size=20, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 5.1, 5.1, 1.3, [
        ("Required пропускная способность: 7,8 миллиарда миль/год", {"size": 12, "italic": True, "color": DEEP}),
        ("Aurora 10-100 грузовиков к 2027 = 1-10M миль пропускная способность", {"size": 12, "italic": True, "color": DEEP}),
        ("AV покрывает 0,01-0,13% дефицита", {"size": 13, "bold": True, "color": RED_WARN}),
    ], line_spacing=1.3)
    # Right: реальный solutions
    rounded_box(slide, 6.3, 1.9, 6.5, 4.5)
    text_box(slide, 6.5, 2.05, 6.1, 0.5, "Что реально решает (не AI)",
             size=16, bold=True, color=MID)
    solutions = [
        ("Программы виз", "Immigration reform, регуляторный rest hours"),
        ("субсидии на обучение CDL", "Government программы — компенсация $3-7K стоимость"),
        ("Реструктуризация зарплат", "Почасовая (включая погрузку) вместо за-милю"),
        ("Качество оборудования", "Newer грузовиков · ~70% оборот вызвано pay+equipment"),
        ("Проектирование рабочих мест", "Local локальный (30-40% оборот) vs дальнобойный (90%+)"),
    ]
    py = 2.65
    for name, desc in solutions:
        text_box(slide, 6.5, py, 6.1, 0.35, "·  " + name,
                 size=12, bold=True, color=DEEP)
        text_box(slide, 6.7, py + 0.35, 5.9, 0.35, desc,
                 size=10, italic=True, color=SLATE)
        py += 0.75
    footer(slide, "Bob Costello (ATA главный экономист): «для всех  неправильных причин» — не AI снизил, спад")
    add_notes(slide, "См. slides/s20-trucker-shortage-false-framing.md speaker notes.")


def s21_highway_failure_matrix(p):
    """s21 — провал matrix Раздел 2."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Четыре причины провала AV-грузоперевозки стартап'ов:\nкапитал / регуляторный / SPAC-bust / симуляция-к-реальности.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Matrix
    rows = [
        ("Капиталоёмкость\nбез выручка", "Argo AI $7B → Oct 2022\nCruise $10B → Dec 2024",
         "Проверять удельная экономика: капитал / коммерческая выручка. Если выручка <5% от капитал — флаг."),
        ("Регуляторная\nнеопределённость", "TuSimple (US-China напряжение Jan 2024)\nTesla NHTSA EA22002",
         "Запрашивать legal team опыт. NHTSA SGO авария reports? Pipeline EA-расследованиеs?"),
        ("SPAC IPO пузырь\n2021-2022", "Embark 16 месяцев SPAC→банкротство\nTuSimple early days",
         "Не вкладывать в до получения выручки SPAC-слияния без 5+ лет коммерческой истории. Публичный надзор убьёт ажиотаж."),
        ("симуляция-к-реальности разрыв", "Starsky (Mar 2020) — general причина\nкаждого недо-выжившего",
         "Запрашивать ratio км в симуляции / км на публичный дороги. Если только sim — серьёзный красный флаг."),
    ]
    cols = ["Причина", "Пример (компания → год)", "Что делать инженеру"]
    col_w = [3.0, 4.0, 5.2]
    col_x = [0.5, 0.5 + 3.1, 0.5 + 3.1 + 4.1]
    y = 1.9
    for i, c in enumerate(cols):
        rectangle(slide, col_x[i], y, col_w[i], 0.5, fill=MID)
        text_box(slide, col_x[i]+0.1, y+0.05, col_w[i]-0.2, 0.4, c,
                 size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.55
    row_h = 1.15
    accent_colors = [RED_WARN, MID, GOLD, TEAL]
    for i, (a, b, c) in enumerate(rows):
        ac = accent_colors[i]
        for j, val in enumerate([a, b, c]):
            rounded_box(slide, col_x[j], y, col_w[j], row_h,
                       stroke=ac if j == 0 else LIGHT, stroke_w=2 if j == 0 else 1)
            text_box(slide, col_x[j]+0.15, y+0.08, col_w[j]-0.3, row_h-0.15, val,
                     size=11 if j > 0 else 12, bold=(j == 0), color=DEEP, line_spacing=1.25)
        y += row_h + 0.05
    footer(slide, "Выжившие применяют поэтапный ввод (ползком-шагом-бегом) против всех 4 причин (см. s19)")
    add_notes(slide, "См. slides/s21-highway-failure-matrix.md speaker notes.")


def s22_starsky_quote(p):
    """s22 — Stefan Seltz-Axmacher quote."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.0,
             "Старски — первый, кто публично признал разрыв.",
             size=26, bold=True, color=DEEP)
    # Large quote box
    rounded_box(slide, 1.0, 1.8, 11.3, 2.8, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    multiline_box(slide, 1.5, 2.0, 10.3, 2.4, [
        ("«Supervised machine learning doesn't live up to the hype.»",
         {"size": 20, "bold": True, "italic": True, "color": DEEP}),
        ("«Контролируемое машинное обучение не оправдывает ажиотажа.»",
         {"size": 14, "italic": True, "color": MID}),
        ("", {"size": 6}),
        ("«Sim-to-real has very real limits.»",
         {"size": 20, "bold": True, "italic": True, "color": DEEP}),
        ("«У разрыва симуляция-к-реальности есть очень реальные пределы.»",
         {"size": 14, "italic": True, "color": MID}),
        ("", {"size": 8}),
        ("— Стефан Зельц-Аксмахер, основатель и CEO Starsky Robotics,",
         {"size": 12, "italic": True, "color": SLATE}),
        ("Medium-эссе «The end of Starsky Robotics», март 2020",
         {"size": 12, "italic": True, "color": SLATE}),
    ], line_spacing=1.25, align=PP_ALIGN.CENTER)
    # Three lessons
    lessons = [
        ("ML не оправдывает ожиданий", "Каждый новый краевой случай требует размеченных данных. До получения выручки масштаба денег не хватает."),
        ("Разрыв симуляция-к-реальности — реальный", "Симуляция не покрывает длинный хвост. ML, обученная только на sim, не обобщается на дороги общего пользования."),
        ("Крупные контракты не материализуются", "Грузоперевозчики — консервативные клиенты. Неиспытанный стек не покупают."),
    ]
    py = 4.8
    for title, desc in lessons:
        rounded_box(slide, 0.5, py, 4.1, 1.7, stroke=LIGHT, stroke_w=1.5)
        text_box(slide, 0.7, py + 0.1, 3.7, 0.5, title,
                 size=12, bold=True, color=MID, line_spacing=1.15)
        text_box(slide, 0.7, py + 0.6, 3.7, 1.0, desc,
                 size=10, color=DEEP, line_spacing=1.3)
        py = py  # next col
    # 3 columns side-by-side
    px = 0.5
    py = 4.8
    for title, desc in lessons:
        rounded_box(slide, px, py, 4.1, 1.7, stroke=LIGHT, stroke_w=1.5)
        text_box(slide, px + 0.2, py + 0.1, 3.7, 0.5, title,
                 size=12, bold=True, color=MID, line_spacing=1.15)
        text_box(slide, px + 0.2, py + 0.6, 3.7, 1.0, desc,
                 size=10, color=DEEP, line_spacing=1.3)
        px += 4.2
    footer(slide, "Жертва первой волны · откровенный разбор полётов от основателя · читать эссе полностью")
    add_notes(slide, "См. slides/s22-starsky-sim-to-real.md speaker notes.")


# ========== SECTION 3 ==========

def s23_section3(p):
    section_divider(p, 3, "Город + последняя миля",
                    ["Robotaxi выжившие: Waymo · Apollo Go · Pony.ai · WeRide",
                     "Pony.ai удельная экономика — Гуанчжоу Nov 2025 (1st) · Шэньчжэнь Feb 2026 (2nd)",
                     "Tesla Robotaxi Austin: только-зрение без HD-map, 14 ДТП за 8 мес.",
                     "Последняя миля: Starship · Coco · Zipline · Nuro разворот",
                     "Cruise GM выход — центральный кейс ($10B → 0)",
                     "Uber Tempe 2018 · Tesla Autopilot 54 смертельных случаев NHTSA"],
                    section_idx=2)


def s24_waymo(p):
    """s24 — Waymo канонический survivor."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Waymo март 2026: 500 000 поездок/неделю, 3 067 машин 5-го поколения.\nHD-карта + LiDAR + удалённый опс + формальный безопасность кейс.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # Hero left
    img_path = ASSETS / "screenshots" / "s24-waymo-self-driving-side.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.5, 4.3)
    attribution(slide, "Waymo Jaguar I-Pace · Wikimedia · CC-BY-SA", x=0.5, y=6.25, w=6.5)
    # Right: stats
    rounded_box(slide, 7.3, 1.9, 5.5, 4.3)
    multiline_box(slide, 7.5, 2.05, 5.1, 4.1, [
        ("Цифры март 2026", {"size": 14, "bold": True, "color": MID}),
        ("500 000 платных поездок/неделю", {"size": 12, "color": DEEP}),
        ("3 067 машин 5-го поколения (NHTSA Dec 2025)", {"size": 12, "color": DEEP}),
        ("10+ городов: Phoenix · SF · LA · Austin · Atlanta · Miami · Dallas · Houston · San Antonio · Orlando", {"size": 11, "color": DEEP}),
        ("14M совокупно поездок за 2025", {"size": 12, "color": DEEP}),
        ("Рост ×10 за 19 месяцев (с 50K/неделю мая 2024)", {"size": 12, "color": GOLD, "bold": True}),
        ("", {"size": 6}),
        ("Стек подход «полный набор сенсоров»", {"size": 14, "bold": True, "color": MID}),
        ("HD-карта · LiDAR · камеры · радар", {"size": 11, "color": DEEP}),
        ("Удалённые операторы (assistance, не driving)", {"size": 11, "color": DEEP}),
        ("Формальный безопасность кейс (регуляторный document)", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Что НЕ публикует", {"size": 14, "bold": True, "color": RED_WARN}),
        ("Per-trip удельная экономика — большая черная коробка на 2026", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Паттерн выживания: поэтапный ввод (ползком-шагом-бегом) · узкий ODD · не переобещают · терпеливый капитал Alphabet")
    add_notes(slide, "См. slides/s24-waymo-survivor.md speaker notes.")


def s25_china_robotaxi(p):
    """s25 — China robotaxi 3 players."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Robotaxi Китая: Apollo Go + Pony.ai + WeRide.\nГосподдержка + единая регуляторика + масштаб.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Three cards horizontal
    players = [
        ("Apollo Go", "Baidu", "240M автономных км глобально · 17M+ заказов · 22 города", LIGHT),
        ("Pony.ai", "NYSE", "300 машин Gen-7 BAIC к Oct 2025 · цель 1000 к концу 2025", MID),
        ("WeRide", "NASDAQ", "Q3 2025 robotaxi выручка 35,3M юаней · +761% YoY", TEAL),
    ]
    card_w = 4.05
    px = 0.5
    py = 1.9
    for name, exchange, desc, color in players:
        rounded_box(slide, px, py, card_w, 3.6, stroke=color, stroke_w=2)
        # Header
        rectangle(slide, px, py, card_w, 1.0, fill=color)
        text_box(slide, px + 0.2, py + 0.1, card_w - 0.4, 0.5, name,
                 size=22, bold=True, color=WHITE)
        text_box(slide, px + 0.2, py + 0.6, card_w - 0.4, 0.35, exchange,
                 size=11, italic=True, color=WHITE)
        # Body
        text_box(slide, px + 0.2, py + 1.2, card_w - 0.4, 2.3, desc,
                 size=13, color=DEEP, line_spacing=1.4)
        px += card_w + 0.1
    # Hero image (Apollo Go) bottom
    img_path = ASSETS / "screenshots" / "s25-apollo-go-rt6.jpg"
    add_image_aspect(slide, img_path, 0.5, 5.7, 6.0, 1.4)
    attribution(slide, "Apollo Go Apollo RT6 · Hubei 2025 · Wikimedia · CC-BY-SA",
                x=0.5, y=7.15, w=6.0)
    # Goldman quote
    rounded_box(slide, 6.8, 5.7, 6.0, 1.4, fill=GOLD_TINT, stroke=GOLD)
    multiline_box(slide, 7.0, 5.8, 5.6, 1.2, [
        ("Goldman Sachs прогноз 2025", {"size": 12, "bold": True, "color": GOLD}),
        ("China robotaxi: ×700 рост от 2025 к $47B к 2035", {"size": 11, "color": DEEP}),
        ("Может обогнать US по числу trips к 2030", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Государственная координация + единая регуляторика на национальном уровне (vs US состояние мозаичный)")
    add_notes(slide, "См. slides/s25-china-robotaxi.md speaker notes.")


def s26_pony_unit_economics(p):
    """s26 — Pony.ai удельная экономика with correct attribution."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Pony.ai: первая положительная unit-экономика на машину.\nГуанчжоу Nov 2025 (1st) · Шэньчжэнь Feb 2026 (2nd).",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image left
    img_path = ASSETS / "screenshots" / "s26-pony-ai-lexus.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 5.5, 4.3)
    attribution(slide, "Pony.ai Lexus RX450h · Guangdong, 2024 · Wikimedia · CC-BY-SA",
                x=0.5, y=6.25, w=5.5)
    # Хронология + key distinction right
    rounded_box(slide, 6.3, 1.9, 6.5, 4.7)
    multiline_box(slide, 6.5, 2.05, 6.1, 4.5, [
        ("Хронология", {"size": 14, "bold": True, "color": MID}),
        ("Ноябрь 2025 · ГУАНЧЖОУ", {"size": 13, "bold": True, "color": GOLD}),
        ("Первый город с положительной операционный прибыль на машину", {"size": 11, "color": DEEP}),
        ("Февраль 2026 · ШЭНЬЧЖЭНЬ", {"size": 13, "bold": True, "color": GOLD}),
        ("Второй город. Расширение паттерн.", {"size": 11, "color": DEEP}),
        ("(SEC 6-K: ~338 юаней daily net income/автомобиль · ~23 заказа/день)", {"size": 10, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("Что означает (и что НЕ означает)", {"size": 14, "bold": True, "color": TEAL}),
        ("Означает: удельная экономика на парк в отдельных городах = положительная.", {"size": 11, "color": DEEP}),
        ("НЕ означает: что компания в целом прибыльна. R&D, маркетинг не покрыты.", {"size": 11, "color": DEEP}),
        ("НЕ означает: что весь сегмент robotaxi прибылен.", {"size": 11, "color": DEEP}),
        ("", {"size": 8}),
        ("Педагогически point", {"size": 14, "bold": True, "color": RED_WARN}),
        ("Первый сигнал, что robotaxi может выйти из venture-жжения. Не magic pill.",
         {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Контраст с Waymo: не публикует за-поездку экономика · большой чёрный ящик на 2026")
    add_notes(slide, "См. slides/s26-pony-ai-unit-economics.md speaker notes.")


def s27_tesla_austin(p):
    """s27 — Tesla Robotaxi Austin current state."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Tesla Robotaxi Austin: 22 июня 2025 — старт.\n~800 000 миль · 14 ДТП за 8 месяцев · только-зрение без HD-карты.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # Hero
    img_path = ASSETS / "screenshots" / "s27-tesla-model-y-myle-2025.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 5.5, 3.2)
    attribution(slide, "Tesla Model Y (2025) · MYLE Festival · Wikimedia · CC-BY-SA",
                x=0.5, y=5.15, w=5.5)
    # Хронология mini below
    rounded_box(slide, 0.5, 5.5, 5.5, 1.5)
    multiline_box(slide, 0.7, 5.6, 5.1, 1.4, [
        ("Хронология 2025-2026", {"size": 12, "bold": True, "color": MID}),
        ("22 июня 2025 · старт с ~10 машин + безопасность monitor", {"size": 10, "color": DEEP}),
        ("Декабрь 2025 · тесты без операторs (employees)", {"size": 10, "color": DEEP}),
        ("Январь 2026 · публичный unsupervised режим", {"size": 10, "color": DEEP}),
        ("Февраль 2026 · 14 ДТП Austin зарегистрировано", {"size": 10, "bold": True, "color": RED_WARN}),
        ("Март-апрель 2026 · downtown Austin · Houston · Dallas", {"size": 10, "color": DEEP}),
    ], line_spacing=1.3)
    # Right: comparison + lesson
    rounded_box(slide, 6.3, 1.9, 6.5, 5.1)
    multiline_box(slide, 6.5, 2.05, 6.1, 5.0, [
        ("Comparison с Waymo (без demonization)", {"size": 14, "bold": True, "color": MID}),
        ("Tesla Austin ~800K миль · Waymo ~10M миль/неделю", {"size": 11, "color": DEEP}),
        ("Объём в 1000 раз меньше — sample size слишком мал", {"size": 11, "italic": True, "color": DEEP}),
        ("«Safer than Waymo» — НЕ доказано · «Хуже Waymo» — НЕ доказано", {"size": 11, "color": DEEP}),
        ("Нужны 3 года данных + сравнение пробег", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 6}),
        ("Только-зрение без HD-карты", {"size": 14, "bold": True, "color": GOLD}),
        ("Tesla верит: сквозной только-зрение стек достаточен для L4", {"size": 11, "color": DEEP}),
        ("Waymo: multi-modal сенсор стек + HD-map (проверенный паттерн)", {"size": 11, "color": DEEP}),
        ("На 2026 — статистически не доказано, кто прав", {"size": 11, "italic": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Урок для инженера", {"size": 14, "bold": True, "color": TEAL}),
        ("Человеческая базовая линия: ~1 смертельный на 100M миль", {"size": 11, "color": DEEP}),
        ("Tesla 14 на 800K — пересчитайте на 100M, сравните с базовая линия", {"size": 11, "color": DEEP}),
        ("Без denominator + базовая линия нельзя оценить «safer/менее safe»", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "«We Robotaxi everywhere by конец of 2026» — Musk Oct 2024. На середину 2026: Austin+Houston+Dallas.")
    add_notes(slide, "См. slides/s27-tesla-robotaxi-austin.md speaker notes.")


def s28_last_mile(p):
    """s28 — Последняя миля 4 cases + Nuro разворот."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Последняя миля: Starship 9M+ доставок · Coco 1000+ роботов LA.\nZipline 100M миль Africa · Nuro разворот 2024 — выход B2C.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # 4 quadrant cards
    items = [
        ("Starship Technologies", "9M+ доставок\n2700+ роботов\n60+ кампусов US",
         "Где работает: кампусы, узкий controlled\nГде не: плотный городской (снег, вандализм)", LIGHT,
         "s28b-starship-tartu-2017.jpg",
         "Starship робот · Tartu, 2017 · Wikimedia"),
        ("Coco Robotics", "1000+ роботов в LA\n500K+ доставок\nЦель 10K единиц",
         "LA + Dallas + Miami + Helsinki + Chicago", MID, None, None),
        ("Zipline (дроны Africa)", "100M миль (Mar 2025)\n2M доставок (Jan 2026)\n22M доз вакцин",
         "$7,6B оценка · $150M State Dept Nov 2025\nProven массовый-развёртывание медицинский Africa", TEAL,
         "s28-zipline-drone-launch.jpg",
         "Zipline Дрон Launch · Wikimedia"),
        ("Nuro разворот 2024", "Exit B2C доставка 2024\nPivot к лицензированию\nавтономный-стека OEM",
         "Урок: B2C доставка не оказался\nприбыльный в American городской context", GOLD, None, None),
    ]
    panel_w = 6.2
    panel_h = 2.6
    for i, (name, stats, where, color, img, attr) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * (panel_w + 0.1)
        y = 1.85 + row * (panel_h + 0.1)
        rounded_box(slide, x, y, panel_w, panel_h, stroke=color, stroke_w=2)
        text_box(slide, x + 0.15, y + 0.1, panel_w - 0.3, 0.45, name,
                 size=14, bold=True, color=color)
        text_box(slide, x + 0.15, y + 0.55, 2.7, 1.0, stats,
                 size=11, color=DEEP, line_spacing=1.3)
        text_box(slide, x + 3.0, y + 0.55, panel_w - 3.2, panel_h - 0.7, where,
                 size=10, italic=True, color=DEEP, line_spacing=1.3)
        # Optional image
        if img:
            img_path = ASSETS / "screenshots" / img
            add_image_aspect(slide, img_path, x + 0.15, y + 1.7, 2.7, 0.85)
            if attr:
                text_box(slide, x + 0.15, y + panel_h - 0.25, 2.7, 0.2, attr,
                         size=8, italic=True, color=SLATE)
    footer(slide, "Тротуарный роботs — узкие ниши · дрон городской US — FAA заблокировано · медицинский Africa — единственный проверенный массовый-развёртывание")
    add_notes(slide, "См. slides/s28-последняя миля-deliveries.md speaker notes.")


def s29_cruise_centerpiece(p):
    """s29 — Cruise GM выход — центральный кейс."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Cruise: 10 миллиардов сожжено, менее 500 миллионов выручки.\nОт 2 окт 2023 (инцидент SF) до 11 дек 2024 (GM выход) — 14 мес. краха.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image
    img_path = ASSETS / "screenshots" / "s29-cruise-bolt-sf.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.5, 3.0)
    attribution(slide, "Cruise Automation Bolt EV in San Francisco · Wikimedia · CC-BY-SA",
                x=0.5, y=4.95, w=6.5)
    # Хронология below
    rounded_box(slide, 0.5, 5.3, 6.5, 1.8)
    multiline_box(slide, 0.7, 5.4, 6.1, 1.7, [
        ("Хронология 2016 → дек 2024", {"size": 12, "bold": True, "color": MID}),
        ("2016 — GM покупает Cruise за $1B+", {"size": 10, "color": DEEP}),
        ("2018-2022 — обширный финансирование, быстрое ODD расширение", {"size": 10, "color": DEEP}),
        ("2 окт 2023 — инцидент SF (волочения инцидент 20 футов)", {"size": 10, "bold": True, "color": RED_WARN}),
        ("24 окт 2023 — California DMV отозвал лицензию", {"size": 10, "bold": True, "color": RED_WARN}),
        ("Late 2023 — mass layoffs, freezing operations", {"size": 10, "color": DEEP}),
        ("11 дек 2024 — GM полный выход · $10B+ убытки · <$500M выручка", {"size": 10, "bold": True, "color": RED_WARN}),
    ], line_spacing=1.2)
    # 4-уровень провал паттерн right
    rounded_box(slide, 7.3, 1.9, 5.5, 5.2)
    multiline_box(slide, 7.5, 2.05, 5.1, 5.0, [
        ("4 уровня провал паттерн", {"size": 14, "bold": True, "color": GOLD}),
        ("", {"size": 4}),
        ("1. Technical", {"size": 12, "bold": True, "color": MID}),
        ("Observe сработал · Decide провалился (тянуть over → катастрофа). Provал 2-й стадии OODA.",
         {"size": 10, "color": DEEP}),
        ("", {"size": 4}),
        ("2. Бизнес режимl", {"size": 12, "bold": True, "color": MID}),
        ("Капитал интенсивность без выручка. $10B vs <$500M = 20:1.", {"size": 10, "color": DEEP}),
        ("", {"size": 4}),
        ("3. Регуляторy / доверие", {"size": 12, "bold": True, "color": MID}),
        ("Cruise представила инцидент менее severely чем был. DMV отозвал не за инцидент, за сокрытие.",
         {"size": 10, "color": DEEP}),
        ("", {"size": 4}),
        ("4. Cultural / organizational", {"size": 12, "bold": True, "color": MID}),
        ("GM-Cruise hybrid культура · misalignment затрудняло поэтапный ввод (ползком-шагом-бегом) discipline.",
         {"size": 10, "color": DEEP}),
        ("", {"size": 4}),
        ("Урок", {"size": 12, "bold": True, "color": RED_WARN}),
        ("Один волочения инцидент + DMV доверие нарушение = убийство program.",
         {"size": 10, "italic": True, "color": DEEP}),
    ], line_spacing=1.2)
    footer(slide, "Паттерн выживания Waymo: поэтапный ввод (ползком-шагом-бегом), conservative ODD, формальный безопасность кейс · Cruise: быстрое ODD ради IPO хронология")
    add_notes(slide, "См. slides/s29-cruise-centerpiece.md speaker notes.")


def s30_uber_tempe(p):
    """s30 — Uber Tempe 2018."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "18 марта 2018, Tempe, Arizona. Элейн Хёрцберг — first AV-пешеход смертельный случай.\nUber отключил заводское AEB, водитель смотрел телевизор.",
             size=18, bold=True, color=DEEP, line_spacing=1.1)
    # Hero image
    img_path = ASSETS / "screenshots" / "s30-uber-volvo-sf.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.5, 3.0)
    attribution(slide, "Uber Self Driving Volvo · San Francisco · Wikimedia · CC-BY-SA",
                x=0.5, y=4.95, w=6.5)
    # Facts below
    rounded_box(slide, 0.5, 5.3, 6.5, 1.8)
    multiline_box(slide, 0.7, 5.4, 6.1, 1.7, [
        ("Что технически произошло", {"size": 12, "bold": True, "color": MID}),
        ("Камера обнаружила пешехода 5,6 сек до удара", {"size": 10, "color": DEEP}),
        ("Классификатор восприятия провалился: не классифицировал как пешеход", {"size": 10, "color": DEEP}),
        ("(Hertzberg была вне пешеходный переход + с велосипедом → вне-распределения)", {"size": 10, "italic": True, "color": SLATE}),
        ("Uber отключил заводское автоматический экстренное торможение (AEB) на Volvo XC90", {"size": 10, "color": RED_WARN}),
        ("Резерв водитель: watching TV (Hulu (стриминг)) во время инцидента", {"size": 10, "color": DEEP}),
    ], line_spacing=1.2)
    # NTSB quotes + lessons right
    rounded_box(slide, 7.3, 1.9, 5.5, 5.2, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    multiline_box(slide, 7.5, 2.05, 5.1, 5.0, [
        ("NTSB quote (HAR-19/03)", {"size": 12, "bold": True, "color": GOLD}),
        ("«Uber ATG's deactivation of its автоматический экстренное торможение система increased  risks associated with testing автоматизированный автомобильs on публичный дороги.»",
         {"size": 10, "italic": True, "color": DEEP}),
        ("", {"size": 4}),
        ("«Uber's недостаточный безопасность культура and недостаточный безопасность risk assessment procedures were cited as factors.»",
         {"size": 10, "italic": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Четыре урока", {"size": 12, "bold": True, "color": MID}),
        ("• ODD критический: обучение данные смещение на пешеходы вне пешеходный переход", {"size": 10, "color": DEEP}),
        ("• Никогда не отключать заводское безопасность systems", {"size": 10, "color": DEEP}),
        ("• Безопасность внимание водителя не reliable (10-15 мин)", {"size": 10, "color": DEEP}),
        ("• Безопасная культура organization имеет значение", {"size": 10, "color": DEEP}),
        ("", {"size": 6}),
        ("Резерв водитель: pleaded guilty endangerment 2023, 3 года supervised испытательный срок",
         {"size": 10, "italic": True, "color": SLATE}),
    ], line_spacing=1.25)
    footer(slide, "Foundation для NHTSA постановление NHTSA on отчётность об авариях · институциональный legacy")
    add_notes(slide, "См. slides/s30-uber-tempe-2018.md speaker notes.")


def s31_tesla_nhtsa(p):
    """s31 — Tesla Autopilot NHTSA 54 смертельных случаев."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "NHTSA SGO Oct 2025: 65 сообщений, 54 проверенный смертельных случаев Tesla Autopilot.\nEA22002: 13 смертельный ДТП с предвидимое неправильное использование.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # Big цифры
    rounded_box(slide, 0.5, 1.9, 4.0, 4.5, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    text_box(slide, 0.7, 2.0, 3.6, 0.4, "NHTSA SGO Oct 2025",
             size=12, italic=True, color=LIGHT)
    text_box(slide, 0.7, 2.5, 3.6, 1.5, "54",
             size=110, bold=True, color=RED_WARN, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(slide, 0.7, 4.2, 3.6, 0.5, "проверенный смертельных случаев",
             size=18, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, 0.7, 4.8, 3.6, 0.4, "(65 reported)",
             size=12, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 5.4, 3.6, 0.95, [
        ("EA22002:", {"size": 12, "bold": True, "color": MID}),
        ("13 смертельный ДТП с предвидимое неправильное использование паттерн", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    # Расследование + lessons right
    rounded_box(slide, 4.8, 1.9, 8.0, 5.0)
    multiline_box(slide, 5.0, 2.05, 7.6, 4.8, [
        ("EA22002 расследование", {"size": 14, "bold": True, "color": MID}),
        ("NHTSA открыл 2022 — для Tesla Autopilot специфически", {"size": 11, "color": DEEP}),
        ("Главный вопрос: достаточно ли мониторинг водителя? warning system?", {"size": 11, "color": DEEP}),
        ("Pattern: водители используют Autopilot для нежных условий (sleep, read)", {"size": 11, "color": DEEP}),
        ("Foreseeable неправильное использование — концепт инженерное standards. 13 смертельный ДТП = структурная проблема дизайна, не вина водителей.",
         {"size": 11, "italic": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Четыре урока", {"size": 14, "bold": True, "color": GOLD}),
        ("• Наименование имеет значение — «Autopilot» / «FSD («полное» самовождение)» приглашение избыточное доверие", {"size": 11, "color": DEEP}),
        ("• Мониторинг водителя обязателен (real вовлечённость проверка, не камера-смотрит-на-лицо)", {"size": 11, "color": DEEP}),
        ("• Краевой случайs в perception (sun засветка, припаркованный экстренное автомобильs)", {"size": 11, "color": DEEP}),
        ("• Только-зрение без HD-map — исследование-stage для L4", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Статистический comparison", {"size": 14, "bold": True, "color": RED_WARN}),
        ("Человеческая базовая линия: ~1 смертельный на 100M миль", {"size": 11, "color": DEEP}),
        ("Tesla 54 смертельных случаев на ~70M миль Autopilot — extrapolation НЕ безопаснее человек", {"size": 11, "italic": True, "color": DEEP}),
        ("Вендор PR утверждает обратное · NHTSA раскрытие показывает структурный паттерн",
         {"size": 11, "italic": True, "color": SLATE}),
    ], line_spacing=1.25)
    footer(slide, "Tesla Автомобиль Безопасность Report использует Tesla's own определение vs human-driver ALL мили · apples-to-oranges")
    add_notes(slide, "См. slides/s31-tesla-autopilot-nhtsa.md speaker notes.")


def s32_urban_failure_matrix(p):
    """s32 — провал matrix Раздел 3."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Четыре урока городского AV: ODD дисциплина / мониторинг водителя /\nнаименование имеет значение / аппаратура ≠ платформа.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Matrix
    rows = [
        ("ODD дисциплина критична", "Cruise (волочения + быстрое расширение)\nUber Tempe (обучение смещение)",
         "При оценке AV-программы: какой ODD? условия за пределами ODD? процесс validation для новый расширение?"),
        ("Мониторинг водителя обязателен", "Tesla Autopilot 54 смертельных случаев NHTSA\nEA22002 13 смертельный с предвидимое неправильное использование",
         "При L2/L3: реальный вовлечённость проверка, не камера-смотрит-на-лицо. Что система делает при distracted driver?"),
        ("Наименование имеет значение", "Tesla «Autopilot» / «FSD («полное» самовождение)» — приглашение избыточное доверие",
         "Никогда не называть L2 ADAS «пилот» / «автономный» / «самоуправление». Использовать терминологию уровней SAE."),
        ("Аппаратура ≠ платформа", "Cruise (GM аппаратура-OEM культура vs программная платформа требования)",
         "Когда OEM покупает software стартап — cultural интеграция plan? GE Predix (GE) lesson повторяется."),
    ]
    cols = ["Урок", "Демонстрирующий кейс", "Что делать инженеру"]
    col_w = [3.2, 4.0, 5.0]
    col_x = [0.5, 0.5 + 3.3, 0.5 + 3.3 + 4.1]
    y = 1.9
    for i, c in enumerate(cols):
        rectangle(slide, col_x[i], y, col_w[i], 0.5, fill=MID)
        text_box(slide, col_x[i]+0.1, y+0.05, col_w[i]-0.2, 0.4, c,
                 size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.55
    row_h = 1.15
    accent_colors = [LIGHT, MID, GOLD, TEAL]
    for i, (a, b, c) in enumerate(rows):
        ac = accent_colors[i]
        for j, val in enumerate([a, b, c]):
            rounded_box(slide, col_x[j], y, col_w[j], row_h,
                       stroke=ac if j == 0 else LIGHT, stroke_w=2 if j == 0 else 1)
            text_box(slide, col_x[j]+0.15, y+0.08, col_w[j]-0.3, row_h-0.15, val,
                     size=11 if j > 0 else 12, bold=(j == 0), color=DEEP, line_spacing=1.25)
        y += row_h + 0.05
    footer(slide, "Связи: ODD + мониторинг водителя · наименование + ODD · аппаратура ≠ платформа (GE Predix (GE) → Cruise GM)")
    add_notes(slide, "См. slides/s32-городской-провал-matrix.md speaker notes.")


# ========== SECTION 4 ==========

def s33_section4(p):
    section_divider(p, 4, "Чрезв. ситуация + рамка",
                    ["Хуситы в Красном море 2024 — 90% drop за 2 месяца",
                     "Suez Ever Given 2021 — 12% мировой торговля",
                     "COVID 2020 снабжение chain обвал",
                     "Дефицит дальнобойщиков 78K — структурная, не AI",
                     "Рамка решения: 5 критериев AI/не AI",
                     "Alternative инструментарий: OR / EOQ / сценарный / на правилах / HITL"],
                    section_idx=3)


def s34_houthi(p):
    """s34 — Houthi Red Sea 2024."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Хуситы в Красном море, конец 2023. За 2 месяца контейнерный трафик упал на 90%.\nОбласть, в которой ML слеп по определению.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # Hero map left
    img_path = ASSETS / "screenshots" / "s34b-prosperity-guardian-map.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.5, 4.0)
    attribution(slide, "Map of Operation Prosperity Guardian · Wikimedia · CC-BY-SA",
                x=0.5, y=5.95, w=6.5)
    # Right panel
    rounded_box(slide, 7.3, 1.9, 5.5, 5.0)
    multiline_box(slide, 7.5, 2.05, 5.1, 4.8, [
        ("Цифры", {"size": 14, "bold": True, "color": MID}),
        ("Декабрь 2023 — атаки начались", {"size": 11, "color": DEEP}),
        ("К февралю 2024: трафик через Red Sea упал на 90% (US DIA)", {"size": 12, "bold": True, "color": RED_WARN}),
        ("Daily объём: 4M т → 1,7M т (−57,5%)", {"size": 11, "color": DEEP}),
        ("Прежде через Red Sea: ~15% морской торговли + ~30% глобальный контейнер", {"size": 11, "color": DEEP}),
        ("+30% транзит время через Cape of Good Hope (Asia-Europe)", {"size": 11, "bold": True, "color": GOLD}),
        ("−9% эффективный глобальный контейнер пропускная способность (J.P. Morgan)", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Что НЕ работало", {"size": 14, "bold": True, "color": RED_WARN}),
        ("ML спрос прогноз полностью вне распределения", {"size": 11, "color": DEEP}),
        ("Оптимизационные решатели — нет данных о новых транзит times", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Что РАБОТАЛО", {"size": 14, "bold": True, "color": TEAL}),
        ("Диспетчеры-люди в исключение-teams (Maersk, MSC, CMA CGM, Hapag-Lloyd)", {"size": 11, "color": DEEP}),
        ("Сценарное планирование (компании с готовыми сценариями)", {"size": 11, "color": DEEP}),
        ("OR с ручное переопределение (Gurobi/OR-Tools перекалибровано)", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Урок: ML по определению слеп на вне распределения событиях · уровень 5 = НЕ AI")
    add_notes(slide, "См. slides/s34-houthi-red-sea.md speaker notes.")


def s35_ever_given(p):
    """s35 — Suez Ever Given 2021."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Suez Ever Given март 2021: 6 дней блокировки, 12% мировой торговли,\n$9,6 миллиарда удержанных грузов. AI не имела роли.",
             size=20, bold=True, color=DEEP, line_spacing=1.1)
    # Hero left: aerial photo of Ever Given
    img_path = ASSETS / "screenshots" / "s35-ever-given-suez.jpg"
    add_image_aspect(slide, img_path, 0.5, 1.9, 6.5, 4.0)
    attribution(slide, "IMO 9811000 EVER GIVEN · Wikimedia · CC-BY-SA",
                x=0.5, y=5.95, w=6.5)
    # Right panel
    rounded_box(slide, 7.3, 1.9, 5.5, 5.0)
    multiline_box(slide, 7.5, 2.05, 5.1, 4.8, [
        ("Факты", {"size": 14, "bold": True, "color": MID}),
        ("23-29 марта 2021 · 6 дней блокировки", {"size": 11, "color": DEEP}),
        ("Ever Given — 400-метровый ULCV (Ultra Large Контейнер Vessel)", {"size": 11, "color": DEEP}),
        ("~12% мировой торговля проходило через Suez", {"size": 11, "color": DEEP}),
        ("$9,6 миллиарда товары застряло (Bloomberg)", {"size": 12, "bold": True, "color": RED_WARN}),
        ("$400M/час delay стоимость estimates (CNBC)", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Что произошло", {"size": 14, "bold": True, "color": MID}),
        ("Sand storm + strong winds + пилот решениеs", {"size": 11, "color": DEEP}),
        ("Канал 200 м ширины · корабль больше в крайних углах", {"size": 11, "color": DEEP}),
        ("AI не имела роли — физика + pilotage", {"size": 12, "bold": True, "color": GOLD}),
        ("", {"size": 6}),
        ("Что помогло разблокировать", {"size": 14, "bold": True, "color": TEAL}),
        ("Engineering, НЕ AI:", {"size": 11, "bold": True, "color": DEEP}),
        ("Tug boats (10+) · dredging · lunar tide 28-29 марта", {"size": 11, "color": DEEP}),
        ("Совокупно effort 5,5 дней — инженерное teams 24/7", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Урок: physical infrastructure problem требует инженерное решения, не AI")
    add_notes(slide, "См. slides/s35-suez-ever-given.md speaker notes.")


def s36_covid_supply_chain(p):
    """s36 — COVID снабжение chain обвал."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "COVID 2020: точно-в-срок не работает на black-swan.\nЧеловеческое исключение-управление спасло.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # 3 phases horizontal
    phases = [
        ("Март-Апрель 2020", "Initial shock",
         "Lockdowns Wuhan → Italy → Spain → US за 6 недель.\nСпрос на потребитель товары всплеск (туалетная бумага, электроника, товары для дома).\nСпрос на travel collapse.", LIGHT),
        ("2020-2021", "Chaos",
         "Ports congested (Long Beach 109 ships waiting Oct 2021).\nКонтейнер rates всплеск 5-10×.\nточно-в-срок сломалось. дефицит СИЗ. дефицит чипов.", RED_WARN),
        ("2022+", "Recalibration",
         "Buffer запасы · многоисточниковый снабжение · nearshoring (Mexico для US, Turkey/Vietnam для Europe).\n«запас «на случай»» вместо «точно-в-срок».", TEAL),
    ]
    phase_w = 4.05
    px = 0.5
    py = 1.9
    for date, name, desc, color in phases:
        rounded_box(slide, px, py, phase_w, 2.5, stroke=color, stroke_w=2)
        rectangle(slide, px, py, phase_w, 0.7, fill=color)
        text_box(slide, px + 0.15, py + 0.05, phase_w - 0.3, 0.3, date,
                 size=11, italic=True, color=WHITE)
        text_box(slide, px + 0.15, py + 0.3, phase_w - 0.3, 0.4, name,
                 size=16, bold=True, color=WHITE)
        text_box(slide, px + 0.15, py + 0.85, phase_w - 0.3, 1.6, desc,
                 size=11, color=DEEP, line_spacing=1.3)
        px += phase_w + 0.1
    # Bottom: what NOT worked vs what worked
    rounded_box(slide, 0.5, 4.6, 6.2, 2.45, stroke=RED_WARN, stroke_w=2)
    multiline_box(slide, 0.7, 4.7, 5.8, 2.3, [
        ("Что НЕ работало", {"size": 14, "bold": True, "color": RED_WARN}),
        ("ML спрос прогнозing — pre-COVID распределение, точность упала на порядок", {"size": 11, "color": DEEP}),
        ("Оптимизационные решатели для запасы — стационарный спрос допущение сломалось", {"size": 11, "color": DEEP}),
        ("Real-time analytics — данные есть, решениеs от человек суждение", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    rounded_box(slide, 6.8, 4.6, 6.0, 2.45, stroke=TEAL, stroke_w=2)
    multiline_box(slide, 7.0, 4.7, 5.6, 2.3, [
        ("Что РАБОТАЛО", {"size": 14, "bold": True, "color": TEAL}),
        ("Управление исключениями человеком — ручное восстановление цепочки поставок", {"size": 11, "color": DEEP}),
        ("Сценарное планирование — готовых сценариев (post-SARS планы 2003)", {"size": 11, "color": DEEP}),
        ("Диверсифицированные снабжение base (многоисточниковый = resilience)", {"size": 11, "color": DEEP}),
        ("Cash reserves для buy-through scarcity", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    footer(slide, "Педагогически: «ML запасы оптимизация» — pitch для стационарный мировой. В мировой с black-swans — buffer + многоисточниковый + сценарный работает лучше.")
    add_notes(slide, "См. slides/s36-covid-supply-chain.md speaker notes.")


def s37_trucker_shortage_structural(p):
    """s37 — дефицит дальнобойщиков структурный problem."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Дефицит дальнобойщиков — структурная проблема трудовой политики, не AI.\nРешается политикой, переподготовкой, условиями труда.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Left: цифры
    rounded_box(slide, 0.5, 1.9, 6.0, 4.5)
    multiline_box(slide, 0.7, 2.05, 5.6, 4.3, [
        ("Цифры (ATA)", {"size": 14, "bold": True, "color": MID}),
        ("78 000 дефицит в 2022 (пик)", {"size": 12, "bold": True, "color": GOLD}),
        ("~60 000 в 2023 («для всех  неправильных причин» — Bob Costello)", {"size": 11, "color": DEEP}),
        ("1,2 миллиона новых нужны за десятилетие", {"size": 11, "color": DEEP}),
        (">90% годовой оборот at large carriers (дальнобойный)", {"size": 11, "color": DEEP}),
        ("~16% female водители (vs ~50% overall рабочая сила)", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Корневые причины (структурные)", {"size": 14, "bold": True, "color": MID}),
        ("· Старение рабочая сила (mean age 52, retirement wave)", {"size": 11, "color": DEEP}),
        ("· Оплата за милю (not почасовая) — эффективный wage под minimum", {"size": 11, "color": DEEP}),
        ("· Lifestyle — дальнобойный недели от дома, family-incompatible", {"size": 11, "color": DEEP}),
        ("· CDL обучение стоимость $3-7K, benefit unclear первые годы", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    # Right: solutions (NOT AI)
    rounded_box(slide, 6.7, 1.9, 6.1, 4.5)
    multiline_box(slide, 6.9, 2.05, 5.7, 4.3, [
        ("Не-AI решения", {"size": 14, "bold": True, "color": TEAL}),
        ("· Программы виз (H-2B расширение для грузовик водители)", {"size": 11, "color": DEEP}),
        ("· субсидии на обучение CDL (government программы)", {"size": 11, "color": DEEP}),
        ("· Реструктуризация зарплат (почасовая включая погрузку время)", {"size": 11, "color": DEEP}),
        ("· Качество оборудования (newer грузовиков, lower пробег)", {"size": 11, "color": DEEP}),
        ("· Проектирование рабочих мест (local локальный vs дальнобойный)", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Почему AV не решит (отзыв s20)", {"size": 14, "bold": True, "color": GOLD}),
        ("78K deficit = 7,8 миллиарда миль пропускная способность needed", {"size": 11, "color": DEEP}),
        ("Aurora 10-100 грузовиков 2027 = 1-10M миль пропускная способность", {"size": 11, "color": DEEP}),
        ("AV покрывает 0,01-0,13% дефицита через 2027 год", {"size": 12, "bold": True, "color": RED_WARN}),
        ("", {"size": 6}),
        ("Урок", {"size": 14, "bold": True, "color": LIGHT}),
        ("Структурные проблемы требуют структурный решений.", {"size": 11, "italic": True, "color": DEEP}),
        ("AI может помочь на margins, не заменить политика.", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.25)
    footer(slide, "Bob Costello (ATA главный экономист): «для всех  неправильных причин» — спад, не AV")
    add_notes(slide, "См. slides/s37-trucker-shortage-structural.md speaker notes.")


def s38_decision_framework(p):
    """s38 — 5 criteria решение фреймворк."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Рамка решения: 5 критериев AI/не AI в логистике.\nСреда · задача · спрос · безопасность · распределение.",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # 5 criteria as cards
    criteria = [
        ("1", "Среда контролируемая?", "Да → AI applicable\n(склад · port · железная дорога)", "Symbotic, Amazon, KONUX", LIGHT),
        ("2", "Чётко определённая оптимизация\n(TSP, VRP, планирование расписаний)?", "Да → OR > RL/ML\n(Gurobi, CPLEX, OR-Tools)", "UPS ORION — $300-400M/год", MID),
        ("3", "Спрос стационарный?", "Да → EOQ + безопасность запас\n+ ABC > ML", "Аудит: какой % SKU реально требует ML? Часто <20%.", TEAL),
        ("4", "Критично для безопасности с\nрегуляторный аудит?", "Да → на правилах + HITL\n(FDA (рег.), FAA, IMO)", "Чёрный ящик ML не работает в регулируемых отраслях.", GOLD),
        ("5", "Событие в распределении?", "Да → ML scoring\nНет → человек-диспетчер\n+ сценарный планирование", "Хуситы, Suez, COVID — вне распределения = ML слеп.", RED_WARN),
    ]
    card_w = 2.42
    px = 0.5
    py = 1.9
    for num, q, ans, ex, color in criteria:
        rounded_box(slide, px, py, card_w, 4.6, stroke=color, stroke_w=2)
        # Number circle
        circle(slide, px + card_w/2 - 0.35, py + 0.15, 0.7, 0.7, fill=color)
        text_box(slide, px + card_w/2 - 0.35, py + 0.15, 0.7, 0.7, num,
                 size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Question
        text_box(slide, px + 0.1, py + 1.0, card_w - 0.2, 1.0, q,
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
        # Answer
        text_box(slide, px + 0.1, py + 2.1, card_w - 0.2, 1.3, ans,
                 size=11, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.3)
        # Example
        rounded_box(slide, px + 0.15, py + 3.5, card_w - 0.3, 1.0, fill=GOLD_TINT, stroke=GOLD)
        text_box(slide, px + 0.25, py + 3.55, card_w - 0.5, 0.9, ex,
                 size=9, italic=True, color=DEEP, line_spacing=1.25, align=PP_ALIGN.CENTER)
        px += card_w + 0.05
    # Bottom takeaway
    rounded_box(slide, 0.5, 6.7, 12.33, 0.5, fill=GOLD_TINT, stroke=GOLD)
    text_box(slide, 0.7, 6.75, 11.9, 0.4,
             "Это не «всегда AI» или «никогда AI». Рамка разбивает нагрузку на категории + правильный инструмент на каждую.",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, "+ дополняется 7 вопросами вендору (слайд s40)")
    add_notes(slide, "См. slides/s38-decision-framework.md speaker notes.")


def s39_alternative_toolkit(p):
    """s39 — Alternative инструментарий matrix."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Альтернативный инструментарий инженера-логиста: 6 классов инструментов.\nAI — один из шести, не «универсальное решение».",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    # Matrix 6 rows × 4 cols
    cols = ["Инструмент", "Задача", "Вендор / открытый исходный код", "Пример"]
    rows = [
        ("OR (Исследование операций)", "Маршрутизация (TSP, VRP), планирование расписаний", "Gurobi, CPLEX, Google OR-Tools",
         "UPS ORION — $300-400M/год"),
        ("Классические запасы", "Управление запасами при стационарный спрос", "(формулы 1913+)",
         "EOQ · безопасность запас · ABC-анализ для большинства SKU"),
        ("Сценарное планирование", "устойчивость к чёрному лебедю", "Shell-style · McKinsey сценарный services",
         "Maersk post-COVID резервирование планирование"),
        ("правиловое зрение", "Контролируемое-env QC", "OpenCV · HALCON · Cognex (визуальный осмотр)",
         "Bottle инспекция на пивоварне (см. lec-11)"),
        ("гибрид CV + обработка сигналов", "многосенсорный осмотр", "Cognex (визуальный осмотр) ЗрениеPro + radar/ультразвук",
         "Контейнер ущерб инспекция в портах"),
        ("Human-in-the-loop (HITL)", "Exception обработка, accountability", "Инструменты управления процессами — Jira, ServiceNow",
         "Maersk исключение teams для Red Sea rerouting"),
    ]
    col_w = [3.0, 2.5, 3.0, 3.7]
    col_x = [0.5]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + 0.05)
    y = 1.9
    # Header
    for i, c in enumerate(cols):
        rectangle(slide, col_x[i], y, col_w[i], 0.4, fill=MID)
        text_box(slide, col_x[i]+0.1, y+0.05, col_w[i]-0.2, 0.3, c,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.45
    row_h = 0.72
    accent_colors = [LIGHT, MID, TEAL, GOLD, LIGHT, MID]
    for i, row in enumerate(rows):
        ac = accent_colors[i]
        for j, val in enumerate(row):
            rounded_box(slide, col_x[j], y, col_w[j], row_h,
                       stroke=ac if j == 0 else LIGHT, stroke_w=2 if j == 0 else 1)
            text_box(slide, col_x[j]+0.1, y+0.05, col_w[j]-0.2, row_h-0.1, val,
                     size=10 if j > 0 else 11, bold=(j == 0), color=DEEP, line_spacing=1.2)
        y += row_h + 0.04
    # Bottom takeaway
    rounded_box(slide, 0.5, 6.7, 12.33, 0.5, fill=GOLD_TINT, stroke=GOLD)
    text_box(slide, 0.7, 6.75, 11.9, 0.4,
             "Инженер-логист, знающий только AI — incomplete engineer. AI ≠ универсальное решение.",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "См. slides/s39-alternative-toolkit.md speaker notes.")


# ========== SECTION 5 ==========

def s40_qa_vendor_questions(p):
    """s40 — Q&A + 7 вендор questions."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 1.2,
             "Семь вопросов вендору на завтра —\nпрактический инструмент для кармана.",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    # 7 questions in 2 columns
    questions = [
        ("1", "Какое сравнение с базовой линией OR\n(Google OR-Tools, Gurobi, CPLEX)?",
         "UPS ORION: OR > ML для чётко поставленной оптимизации.\nЕсли поставщик не делал сравнения — красный флаг.", LIGHT),
        ("2", "Какой ваш ODD,\nи как валидируется новое расширение?",
         "Cruise dragging инцидент Oct 2023 — провал именно ODD-дисциплины.\nРасширение без обширной валидации — анти-паттерн.", MID),
        ("3", "Какой ваш стек мониторинга водителя\n(для L2/L3)?",
         "Tesla EA22002: 13 смертельных ДТП с предсказуемым неправильным использованием.\nСтруктурная проблема дизайна, не вина водителей.", TEAL),
        ("4", "Отношение км в симуляции / км на дорогах общего пользования?",
         "Starsky разрыв симуляция-к-реальности.\nЕсли только симуляция — серьёзный красный флаг.", GOLD),
        ("5", "Частота ошибок на сезонных сдвигах распределения\n(Чёрная пятница, Рождество)?",
         "Сдвиг распределения на сезонных пиках убивает ML-модели,\nобученные на данных вне пиков.", LIGHT),
        ("6", "Какие сертификации\n(FDA Part 11, ATEX, ISO 26262, NHTSA SGO)?",
         "Регуляторный аудит обязателен в категориях, критичных для безопасности.\nЧёрный ящик ML не проходит аудит.", MID),
        ("7", "Какова удельная экономика\n(на машину / маршрут / тонну)?",
         "Pony.ai первая robotaxi с положительной операционной прибылью на автомобиль\n(Гуанчжоу ноябрь 2025, Шэньчжэнь февраль 2026).", RED_WARN),
    ]
    px = 0.5
    py = 1.7
    card_w = 6.05
    card_h = 1.55
    for i, (num, q, why, color) in enumerate(questions):
        col = i % 2
        row = i // 2
        x = 0.5 + col * (card_w + 0.1)
        y = 1.7 + row * (card_h + 0.05)
        rounded_box(slide, x, y, card_w, card_h, stroke=color, stroke_w=2)
        # Number circle
        circle(slide, x + 0.15, y + 0.2, 0.5, 0.5, fill=color)
        text_box(slide, x + 0.15, y + 0.2, 0.5, 0.5, num,
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Question
        text_box(slide, x + 0.75, y + 0.1, card_w - 0.85, 0.7, q,
                 size=12, bold=True, color=DEEP, line_spacing=1.2)
        # Почему
        text_box(slide, x + 0.15, y + 0.85, card_w - 0.3, 0.65, "Почему: " + why,
                 size=10, italic=True, color=SLATE, line_spacing=1.25)
    footer(slide, "Дополняет 5-критерийную рамку слайда s38 · окупаемость лекции 13")
    add_notes(slide, "См. slides/s40-qa-vendor-questions.md speaker notes.")


def s41_closing_hero_noc(p):
    """s41 — Closing hero NOC bridge to lec-14."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    # Hero image left (≥40% area = ≥40 sq inch on 13.33×7.5 slide; 8.0×5.5 = 44 sq inch)
    img_path = ASSETS / "screenshots" / "s41-noc-iupui.jpg"
    add_image_aspect(slide, img_path, 0.5, 0.5, 8.0, 5.5)
    attribution(slide, "Network Операционные Center · IUPUI · Wikimedia · CC-BY-SA",
                x=0.5, y=6.05, w=8.0)
    # Right: bridge text (narrower column)
    multiline_box(slide, 8.8, 0.6, 4.2, 5.5, [
        ("Мост к Лекции 14", {"size": 14, "bold": True, "color": LIGHT}),
        ("Другая среда,", {"size": 24, "bold": True, "color": DEEP}),
        ("те же 5 вопросов", {"size": 24, "bold": True, "color": DEEP}),
        ("", {"size": 12}),
        ("Следующая лекция —", {"size": 13, "color": MID}),
        ("телекоммуникации,", {"size": 13, "color": MID}),
        ("сетевая инфраструктура,", {"size": 13, "color": MID}),
        ("кибербезопасность.", {"size": 13, "color": MID}),
        ("", {"size": 8}),
        ("· AI помогает SOC-аналитику", {"size": 11, "color": DEEP}),
        ("· OR + правиловая детекция —", {"size": 11, "color": DEEP}),
        ("  mainstream", {"size": 11, "color": DEEP}),
        ("· ML — в-распределении силён,", {"size": 11, "color": DEEP}),
        ("  вне — слеп", {"size": 11, "color": DEEP}),
        ("· HITL обязателен на исключениях", {"size": 11, "color": DEEP}),
        ("", {"size": 8}),
        ("Среда меняется.", {"size": 14, "bold": True, "color": GOLD}),
        ("Критическое суждение — нет.", {"size": 14, "bold": True, "color": GOLD}),
    ], line_spacing=1.15)
    # Recap callout
    rounded_box(slide, 0.5, 6.1, 12.33, 1.0, fill=GOLD_TINT, stroke=GOLD, stroke_w=2)
    multiline_box(slide, 0.7, 6.2, 11.9, 0.9, [
        ("Выжившие уважают среду + остаются в узком ODD + не переобещают.",
         {"size": 14, "bold": True, "color": DEEP, "align": PP_ALIGN.CENTER}),
        ("Cruise vs Waymo — обоим тот же стек. Waymo выжил, потому что cautious в ODD расширение.",
         {"size": 12, "italic": True, "color": DEEP, "align": PP_ALIGN.CENTER}),
    ], align=PP_ALIGN.CENTER, line_spacing=1.25)
    add_notes(slide, "См. slides/s41-closing-hero-noc-bridge.md speaker notes.")


# ========== BUILD ==========

def main():
    p = setup_pres()
    # Section 0
    s01_hero_three_pictures(p)
    s02_cover(p)
    s03_lecture_map(p)
    s04_glossary(p)
    s05_keystone_ladder(p)
    # Section 1
    s06_section1(p)
    s07_symbotic_walmart(p)
    s08_amazon_robotics(p)
    s09_amr_locus(p)
    s10_port_automation(p)
    s11_rail_konux(p)
    s12_discrete_failure_matrix(p)
    # Section 2
    s13_section2(p)
    s14_aurora(p)
    s15_mobileye_kamaz(p)
    s16_ups_orion(p)
    s17_av_bankruptcy_timeline(p)
    s18_cumulative_20b(p)
    s19_survivor_consolidation(p)
    s20_trucker_shortage_false(p)
    s21_highway_failure_matrix(p)
    s22_starsky_quote(p)
    # Section 3
    s23_section3(p)
    s24_waymo(p)
    s25_china_robotaxi(p)
    s26_pony_unit_economics(p)
    s27_tesla_austin(p)
    s28_last_mile(p)
    s29_cruise_centerpiece(p)
    s30_uber_tempe(p)
    s31_tesla_nhtsa(p)
    s32_urban_failure_matrix(p)
    # Section 4
    s33_section4(p)
    s34_houthi(p)
    s35_ever_given(p)
    s36_covid_supply_chain(p)
    s37_trucker_shortage_structural(p)
    s38_decision_framework(p)
    s39_alternative_toolkit(p)
    # Section 5
    s40_qa_vendor_questions(p)
    s41_closing_hero_noc(p)

    p.save(str(OUT))
    print(f"Saved {OUT} ({len(p.slides)} slides)")


if __name__ == "__main__":
    main()
