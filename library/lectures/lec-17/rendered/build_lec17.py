"""
build_lec17.py — Лекция 17 «Систематизация знаний и навыков — инженерная карта AI»
(capstone). REBUILT v2 (Phase 8c-2) под 37-slide структуру (deck.yaml v2).

Source-of-truth: deck.yaml v2 (37 slides) + slides/s01..s37*.md.
Issue #145 · downstream от chapter (book-first).

Palette LOCKED v3: Ocean (#21295C / #065A82 / #1C7293) + Teal (#028090) + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9).

Slide order (v2):
  s01 title cover | s02 keystone | s03 Р1-divider | s04-s09 Р1 | s10 Р2-divider
  s11-s17 Р2 | s18 Р3-divider | s19 lecture-map-16 | s20-s26 Р3 | s27 Р4-divider
  s28-s31 Р4 | s32 cheatsheets-overview | s33-s36 cheatsheet-previews | s37 QA

Scatter PNGs (render_scatter.py): s01-hero / s02-keystone / s20-batch1 /
s21-batch2 / s22-batch3-full / s23-cluster-ur / s24-cluster-ul / s25-cluster-lr /
s26-empty-quadrants / s36-master-poster.

Case images embedded (assets/screenshots): see-and-spray s05 / crowdstrike s06 /
ups-orion s08 / aidoc s13 / waymo s14 / symbotic s15 / klarna s16 / alphafold s23 /
monarch s24 / epic-sepsis s25 / uber-tempe s28 / arup-deepfake s29 / getty-stability s30.

Build:
  python3 render_scatter.py     # scatter PNGs (run first)
  python3 build_lec17.py        # generates lec-17.pptx (notes = stubs)
  python3 inject_notes.py       # injects FULL speaker notes from slides/*.md
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT  = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT  = RGBColor(0xE6, 0xF2, 0xF4)
GREEN_TINT = RGBColor(0xE4, 0xF1, 0xE8)
BLUE_TINT  = RGBColor(0xE6, 0xEE, 0xF5)
SOFT_GREY  = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY  = RGBColor(0x4A, 0x55, 0x6B)
ROADMAP    = RGBColor(0xD9, 0xE2, 0xEC)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
CHARTS = ASSETS / "charts"
SHOTS = ROOT.parent / "assets" / "screenshots"   # library/lectures/lec-17/assets/screenshots
OUT = ROOT / "lec-17.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"

OK = "✓"
WARN = "⚠"
NO = "✗"


# --------------------------------------------------------------------- helpers
def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"
    for el in sppr.findall(ns):
        sppr.remove(el)
    etree.SubElement(sppr, ns)


def text_box(slide, x, y, w, h, text, *, size=16, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.12):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except Exception: pass
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline_box(slide, x, y, w, h, lines, *, size=14, bold=False, color=DEEP,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
                  line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except Exception: pass
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        if isinstance(txt, list):
            for seg, sopts in txt:
                r = p.add_run(); r.text = seg
                r.font.name = sopts.get("font", font)
                r.font.size = Pt(sopts.get("size", size))
                r.font.bold = sopts.get("bold", bold)
                r.font.italic = sopts.get("italic", False)
                r.font.color.rgb = sopts.get("color", color)
        else:
            r = p.add_run(); r.text = txt
            r.font.name = opts.get("font", font)
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", color)
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None, stroke_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_w)
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, w, h, *, fill=MID, stroke=None, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_w)
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, *, fill=LIGHT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def up_arrow(slide, x, y, w, h, *, fill=TEAL):
    shp = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, color=WHITE, size=11, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    tf = shp.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return shp


def add_image_aspect(slide, path, x, y, w, h):
    """Add picture preserving aspect ratio (centered in box)."""
    p = Path(path)
    if not p.exists():
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        text_box(slide, x, y + h / 2 - 0.3, w, 0.6, f"[нет: {p.name}]",
                 size=10, color=SLATE, align=PP_ALIGN.CENTER)
        return None
    try:
        with Image.open(p) as img:
            iw, ih = img.size
        img_ratio = iw / ih
        box_ratio = w / h
        if img_ratio > box_ratio:
            new_w = w; new_h = w / img_ratio
            cx = x; cy = y + (h - new_h) / 2
        else:
            new_h = h; new_w = h * img_ratio
            cx = x + (w - new_w) / 2; cy = y
        return slide.shapes.add_picture(str(p), Inches(cx), Inches(cy),
                                        width=Inches(new_w), height=Inches(new_h))
    except Exception as e:
        print(f"Image fail {path}: {e}")
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        return None


def add_image_crop_fill(slide, path, x, y, w, h):
    """Add picture filling box exactly, cropping overflow (center crop)."""
    p = Path(path)
    if not p.exists():
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        return None
    try:
        with Image.open(p) as img:
            iw, ih = img.size
        img_ratio = iw / ih
        box_ratio = w / h
        pic = slide.shapes.add_picture(str(p), Inches(x), Inches(y),
                                       width=Inches(w), height=Inches(h))
        if img_ratio > box_ratio:
            # image wider -> crop left/right
            crop = (1 - box_ratio / img_ratio) / 2
            pic.crop_left = crop; pic.crop_right = crop
        else:
            crop = (1 - img_ratio / box_ratio) / 2
            pic.crop_top = crop; pic.crop_bottom = crop
        return pic
    except Exception as e:
        print(f"Image crop fail {path}: {e}")
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        return None


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# 5 roadmap sections (no minutes — divider/cover only)
SECTIONS = ["Критерии", "Лестница", "Карта 16", "Провалы", "Карточки"]


def roadmap_bar(slide, current_section):
    bar_y = 6.95; bar_h = 0.32; total_w = 12.33
    seg_w = total_w / 5
    for i, name in enumerate(SECTIONS):
        x = 0.5 + i * seg_w
        active = (i == current_section)
        rectangle(slide, x, bar_y, seg_w - 0.06, bar_h,
                  fill=GOLD if active else ROADMAP)
        text_box(slide, x, bar_y + 0.025, seg_w - 0.06, bar_h - 0.04,
                 f"{i+1}. {name}", size=11, bold=active,
                 color=DEEP if active else SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def assertion_title(slide, text, *, y=0.42, size=24, w=12.33, x=0.5):
    text_box(slide, x, y, w, 1.0, text, size=size, bold=True, color=DEEP,
             line_spacing=1.06)


def footer(slide, text, *, y=7.06, color=LIGHT, size=11.5):
    text_box(slide, 0.5, y, 12.33, 0.38, text, size=size, italic=True,
             color=color, align=PP_ALIGN.LEFT, line_spacing=1.05)


def attribution(slide, text, *, x=0.5, y=7.06, w=12.33, align=PP_ALIGN.LEFT, size=10):
    text_box(slide, x, y, w, 0.3, text, size=size, italic=True, color=SLATE,
             align=align)


def img_attribution(slide, x, y, w, text, *, h=0.34):
    """Small attribution caption strip under an embedded real photo (RU)."""
    text_box(slide, x + 0.06, y, w - 0.12, h, text,
             size=8, italic=True, color=SLATE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)


def gold_callout(slide, x, y, w, h, text, *, size=14, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE):
    rounded_box(slide, x, y, w, h, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.75)
    text_box(slide, x + 0.2, y + 0.06, w - 0.4, h - 0.12, text,
             size=size, bold=True, color=DEEP, anchor=anchor,
             align=align, line_spacing=1.12)


def icon_badge(slide, x, y, d, glyph, *, fill=MID, color=WHITE, size=18):
    circle(slide, x, y, d, d, fill=fill)
    text_box(slide, x, y, d, d, glyph, size=size, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# SECTION 0 — Title + keystone (s01-s02)
# ====================================================================

def s01(p):
    """hero_cover — CLEAN title slide. Title dominant, scatter as supporting
    visual on the right (~45% width). NO roadmap bar on cover."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    # left text column
    text_box(slide, 0.7, 0.7, 6.4, 0.4,
             "ПРИМЕНЕНИЕ ИСКУССТВЕННОГО ИНТЕЛЛЕКТА · ИТОГОВАЯ ЛЕКЦИЯ КУРСА",
             size=12, bold=True, color=LIGHT)
    # capstone marker chip
    chip(slide, 0.7, 1.25, 2.95, 0.45, "CAPSTONE — ФИНАЛ КУРСА",
         fill=GOLD, color=DEEP, size=12.5)
    # Lecture title — dominant
    multiline_box(slide, 0.7, 2.0, 6.5, 2.7, [
        ("Лекция 17", {"size": 22, "bold": True, "color": LIGHT}),
        ("", {"size": 6}),
        ("Систематизация знаний", {"size": 33, "bold": True, "color": DEEP}),
        ("и навыков — инженерная", {"size": 33, "bold": True, "color": DEEP}),
        ("карта AI", {"size": 33, "bold": True, "color": DEEP}),
    ], line_spacing=1.04)
    # hook line in motif box, gold keywords
    rounded_box(slide, 0.7, 5.1, 6.5, 1.45, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    multiline_box(slide, 0.95, 5.32, 6.0, 1.05, [
        ([("Где AI ", {"size": 21, "bold": True, "color": DEEP}),
          ("работает", {"size": 21, "bold": True, "color": GOLD}),
          (", где — ", {"size": 21, "bold": True, "color": DEEP}),
          ("нет", {"size": 21, "bold": True, "color": GOLD}),
          (",", {"size": 21, "bold": True, "color": DEEP})], {}),
        ([("и как это ", {"size": 21, "bold": True, "color": DEEP}),
          ("понять", {"size": 21, "bold": True, "color": GOLD}),
          ("?", {"size": 21, "bold": True, "color": DEEP})], {}),
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    # right supporting scatter (~45% width)
    rounded_box(slide, 7.5, 0.95, 5.35, 5.6, fill=WHITE, stroke=LIGHT, stroke_w=1.5)
    add_image_aspect(slide, CHARTS / "s01-hero-scatter.png", 7.62, 1.07, 5.11, 5.05)
    text_box(slide, 7.62, 6.18, 5.11, 0.32,
             "Карта 16 отраслей курса на одной плоскости",
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)
    attribution(slide, "Карта составлена по 16 отраслевым лекциям курса · оси SAE J3016 (адаптация)",
                x=0.7, y=7.08, w=12.0)
    add_notes(slide, "s01")
    return slide


def s02(p):
    """schema_quadrant — keystone reveal: 2 axes + 4 labeled quadrants, no points."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Карта применимости ИИ: горизонталь — нужен ли AI, вертикаль — сколько ему доверить.",
        y=0.4, size=22)
    rounded_box(slide, 0.5, 1.35, 9.15, 5.05, fill=WHITE, stroke=LIGHT, stroke_w=1.5)
    add_image_aspect(slide, CHARTS / "s02-keystone-quadrants.png", 0.62, 1.48, 8.91, 4.78)
    # Right explanation column
    rx = 9.85
    rounded_box(slide, rx, 1.35, 2.98, 5.05, fill=SURFACE, stroke=LIGHT)
    multiline_box(slide, rx + 0.22, 1.58, 2.55, 3.7, [
        ("Две оси", {"size": 16, "bold": True, "color": MID}),
        ("", {"size": 5}),
        ("Горизонталь — нужен ли AI вообще (применимость).",
         {"size": 12.5, "color": DARK_GREY}),
        ("", {"size": 5}),
        ("Вертикаль — сколько автономии доверить (L0→L5).",
         {"size": 12.5, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("Высокое значение — не «лучше», а «AI подходит» / «AI делает больше».",
         {"size": 12.5, "bold": True, "color": DEEP}),
    ], line_spacing=1.16)
    gold_callout(slide, rx + 0.22, 5.42, 2.55, 0.82,
                 "Верхне-левый квадрант обязан оставаться пустым.", size=12)
    footer(slide, "Несущая ось capstone — предъявлена до первого погружения. Источник осей: SAE J3016 (адаптация).")
    add_notes(slide, "s02")
    return slide


# ====================================================================
# Section dividers (s03 / s10 / s18 / s27)
# ====================================================================

def section_divider(p, number, title, subtitle, tag, section_idx, notes_key):
    slide = blank(p); set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=section_idx)
    # large decorative number
    text_box(slide, 0.3, 0.9, 5.2, 4.6, number, size=300, bold=True, color=SOFT_GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title + subtitle right
    multiline_box(slide, 5.2, 2.1, 7.6, 3.2, [
        (f"Раздел {number}", {"size": 18, "bold": True, "color": GOLD}),
        ("", {"size": 8}),
        (title, {"size": 33, "bold": True, "color": DEEP}),
        ("", {"size": 12}),
        (subtitle, {"size": 15.5, "color": DARK_GREY}),
    ], line_spacing=1.16)
    gold_callout(slide, 0.5, 6.05, 12.33, 0.72, tag, size=13.5)
    add_notes(slide, notes_key)
    return slide


def s03(p):
    return section_divider(p, "1", "Когда AI применять, когда нет",
        "Семь критериев — горизонтальная ось карты как диагностический чек-лист.",
        "7 критериев · последовательные ворота · полностью про границы", 0, "s03")


def s10(p):
    return section_divider(p, "2", "Лестница автономии L0→L5",
        "Вертикальная ось карты — формальная шкала с маппингом локальных шкал курса в единую.",
        "6 ступеней · кросс-отраслевой маппинг · антипаттерн на каждой", 1, "s10")


def s18(p):
    return section_divider(p, "3", "16 отраслей на одной карте",
        "Точечная диаграмма строится слой за слоем + кластерный анализ.",
        "точечная карта · 1 плоскость · ~20 точек · 4 квадранта", 2, "s18")


def s27(p):
    return section_divider(p, "4", "Двенадцать провалов курса. Что выучили?",
        "Каждый провал — повторяющийся паттерн, не ошибка одной компании.",
        "кросс-отраслевые · уроки → альтернативы · полностью про границы", 3, "s27")


# ====================================================================
# SECTION 1 — Семь критериев (s04-s09)
# ====================================================================

def s04(p):
    """assertion_visual — 7 criteria overview (numbered rows w/ glyph badges)."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Семь критериев работают как последовательные ворота: один однозначный «нет» — "
        "и задача не для полной автономии AI.", y=0.4, size=20)
    crit = [
        ("Закрытая петля или открытая среда?", "Есть ли быстрая измеримая обратная связь?", "↻", MID),
        ("Достаточно ли обучающих данных", "и совпадает ли распределение с эксплуатацией?", "▤", MID),
        ("Повторяемость и объём", "задача частая или штучная?", "⟳", MID),
        ("Цена ошибки и радиус разрушения", "сколько стоит ошибка и как широко расходится?", "⚠", GOLD),
        ("Доступность эталона", "есть ли быстрая проверка правильности?", "◉", MID),
        ("Объяснимость и аудит", "нужно ли объяснить каждое решение?", "◎", MID),
        ("Экономика против базовой линии", "окупается ли AI против классической альтернативы?", "⚖", MID),
    ]
    gy = 1.55; rh = 0.58; gap = 0.075
    for i, (head, sub, glyph, col) in enumerate(crit):
        y = gy + i * (rh + gap)
        rounded_box(slide, 0.5, y, 8.5, rh, fill=SURFACE, stroke=LIGHT, stroke_w=1.25)
        icon_badge(slide, 0.62, y + 0.09, rh - 0.18, glyph, fill=col, size=17)
        multiline_box(slide, 1.32, y + 0.06, 7.6, rh - 0.08, [
            ([(f"{i+1}. ", {"size": 13.5, "bold": True, "color": col}),
              (head, {"size": 13.5, "bold": True, "color": DEEP})], {}),
            (sub, {"size": 11, "color": DARK_GREY}),
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    # right rule callout
    rx = 9.25
    rounded_box(slide, rx, 1.55, 3.58, 4.7, fill=BLUE_TINT, stroke=MID, stroke_w=1.5)
    multiline_box(slide, rx + 0.24, 1.82, 3.1, 4.2, [
        ("Правило применения", {"size": 15, "bold": True, "color": MID}),
        ("", {"size": 8}),
        ([("Один ", {"size": 14, "color": DEEP}),
          ("✗ → СТОП", {"size": 15, "bold": True, "color": GOLD}),
          (": полная автономия не подходит.", {"size": 13.5, "color": DEEP})], {}),
        ("", {"size": 8}),
        ("Несколько ⚠ → начинайте с advisory + человек в петле.",
         {"size": 13.5, "color": DARK_GREY}),
        ("", {"size": 8}),
        ("Все ✓ → пилот с явными точками GO / NO-GO.",
         {"size": 13.5, "color": DARK_GREY}),
    ], line_spacing=1.18)
    footer(slide, "Критерии работают в комплексе — пересекаются и усиливают друг друга. Разбираем их парами.")
    add_notes(slide, "s04")
    return slide


def two_col_compare(slide, *, left_title, left_color, left_items, right_title,
                    right_color, right_items, top=1.5, h=4.0):
    """Equal-height two-column comparison. Each item: (head, body)."""
    cw = 5.95; gx = 0.5; gap = 0.43
    cols = [
        (gx, left_title, left_color, left_items, GREEN_TINT),
        (gx + cw + gap, right_title, right_color, right_items, GOLD_TINT),
    ]
    for cx, title, col, items, tint in cols:
        rounded_box(slide, cx, top, cw, h, fill=tint, stroke=col, stroke_w=1.75)
        # header bar
        rounded_box(slide, cx + 0.18, top + 0.16, cw - 0.36, 0.52, fill=col,
                    stroke=None, radius=0.2)
        text_box(slide, cx + 0.3, top + 0.18, cw - 0.6, 0.48, title, size=15,
                 bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # items
        iy = top + 0.86
        n = len(items)
        ih = (h - 1.02) / n
        for head, body in items:
            multiline_box(slide, cx + 0.3, iy + 0.04, cw - 0.6, ih - 0.06, [
                (head, {"size": 13, "bold": True, "color": DEEP}),
                (body, {"size": 11, "color": DARK_GREY}),
            ], line_spacing=1.06, anchor=MSO_ANCHOR.TOP)
            iy += ih


def s05(p):
    """comparison — criteria 1+2; closed-loop vs open-env, with See & Spray photo."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Закрытая петля с быстрой обратной связью — AI работает; открытая среда "
        "с распределённым сдвигом — AI системно проваливается.", y=0.4, size=19)
    two_col_compare(slide,
        left_title="Закрытая петля — AI работает", left_color=TEAL,
        left_items=[
            ("See & Spray — агро (фото)", "контролируемый ряд, −50% гербицидов: ≈1→0,5 фунт/акр на 5 млн акров (0,55% пашни США)"),
            ("Aidoc — медицина (радиология)", "стандартизованные снимки, человек в петле, FDA-режим L1"),
            ("Складская робототехника — логистика", "миллионы операций в контролируемой среде"),
        ],
        right_title="Открытая среда — AI проваливается", right_color=GOLD,
        right_items=[
            ("Zillow — недвижимость", "COVID-сдвиг, $304 млн списания, 25% сокращений (≈2 000/8 000), ноя. 2021"),
            ("Monarch — агро", "открытое поле, погода; 38% сокращений (≈53/140), янв. 2025"),
            ("Cruise — транспорт", "протащил пешехода 20 футов 02.10.2023; GM закрыл подразделение 10.12.2024"),
        ], top=1.45, h=3.55)
    # See & Spray photo strip + gold callout for criterion 2
    iy = 5.2
    rounded_box(slide, 0.5, iy, 4.1, 1.55, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    add_image_aspect(slide, SHOTS / "see-and-spray.jpg", 0.62, iy + 0.1, 3.86, 1.18)
    img_attribution(slide, 0.62, iy + 1.25, 3.86, "John Deere · Wikimedia · CC BY 2.0")
    gold_callout(slide, 4.85, iy, 7.98, 1.55,
        "Критерий 2: даже закрытая среда не спасёт без данных, отражающих эксплуатацию. "
        "Epic Sepsis — вендор заявлял AUC 0,76, внешняя валидация на 38 000 пациентов дала 0,63.",
        size=13)
    add_notes(slide, "s05")
    return slide


def s06(p):
    """comparison — criteria 3+4 blast-radius table + CrowdStrike photo."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "AI окупается на повторяемых задачах с низкой ценой ошибки — и опасен там, "
        "где цена ошибки умножается на большой радиус разрушения.", y=0.4, size=19)
    # Criterion 3 block (left)
    rounded_box(slide, 0.5, 1.5, 7.7, 1.7, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 0.72, 1.62, 7.3, 0.34, "Критерий 3 — повторяемость и объём",
             size=14, bold=True, color=MID)
    multiline_box(slide, 0.72, 2.0, 7.3, 1.1, [
        ([("GitHub Copilot — разработка ПО: ", {"size": 12, "bold": True, "color": DEEP}),
          ("миллиарды дополнений/день; 20+ млн платных (из ≈28 млн разработчиков GitHub); 46% кода. Ошибка — нажать Esc. AI окупается с запасом.",
           {"size": 12, "color": DARK_GREY})], {}),
        ([("Архитектура крупной системы: ", {"size": 12, "bold": True, "color": DEEP}),
          ("раз в год-полтора, огромная цена ошибки. AI не окупается — нужен сеньор.",
           {"size": 12, "color": DARK_GREY})], {}),
    ], line_spacing=1.1)
    # CrowdStrike photo (right of crit3)
    rounded_box(slide, 8.4, 1.5, 4.43, 1.7, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    add_image_aspect(slide, SHOTS / "crowdstrike.jpg", 8.52, 1.6, 4.19, 1.32)
    img_attribution(slide, 8.52, 2.9, 4.19, "BSOD на табло LaGuardia 19.07.2024 · Wikimedia · CC BY-SA 4.0")
    # Criterion 4 blast-radius table
    text_box(slide, 0.5, 3.42, 12.33, 0.34,
             "Критерий 4 — цена ошибки × радиус разрушения (по возрастанию радиуса)",
             size=14, bold=True, color=MID)
    cols_x = [0.5, 4.7, 7.3, 9.6]; cols_w = [4.2, 2.6, 2.3, 3.23]
    headers = ["Применение", "Радиус", "Обнаружение", "Откат"]
    rows = [
        ("Copilot — ошибка автодополнения", "1 разработчик", "секунды", "секунды (Ctrl-Z)", False),
        ("Aidoc — пропуск (в L1)", "1 пациент", "минуты-дни", "радиолог проверит", False),
        ("Waymo — инцидент", "1-несколько", "секунды", "минуты", False),
        ("CrowdStrike Falcon (19.07.2024)", "8,5 млн устройств", "минуты", "часы (ребут каждого)", True),
    ]
    ty = 3.82; rh = 0.58
    # header row
    for j, hdr in enumerate(headers):
        rectangle(slide, cols_x[j], ty, cols_w[j] - 0.06, 0.4, fill=MID)
        text_box(slide, cols_x[j] + 0.08, ty + 0.03, cols_w[j] - 0.2, 0.34, hdr,
                 size=11.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    ry = ty + 0.44
    for app, radius, detect, rollback, gold in rows:
        vals = [app, radius, detect, rollback]
        for j, v in enumerate(vals):
            rounded_box(slide, cols_x[j], ry, cols_w[j] - 0.06, rh,
                        fill=GOLD_TINT if gold else SURFACE,
                        stroke=GOLD if gold else SOFT_GREY,
                        stroke_w=1.75 if gold else 0.75, radius=0.05)
            text_box(slide, cols_x[j] + 0.1, ry + 0.04, cols_w[j] - 0.24, rh - 0.08, v,
                     size=11.5, bold=gold, color=DEEP if not gold else DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        ry += rh + 0.06
    footer(slide, "CrowdStrike отличается на порядки. $5+ млрд ущерба клиентам (оценка Parametrix по Fortune 500 — не убытки самой CrowdStrike). Широкий радиус не запрещает систему — требует строже дисциплины развёртывания.",
           size=10.5)
    add_notes(slide, "s06")
    return slide


def s07(p):
    """comparison — criteria 5+6: Pearl ladder + EU AI Act / Apple Card."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Без быстрого эталона и без объяснимости AI остаётся на advisory — даже без "
        "злого умысла чёрная коробка разрушает доверие.", y=0.4, size=19)
    # LEFT: criterion 5 — Pearl ladder (bottom-aligned)
    rounded_box(slide, 0.5, 1.5, 6.0, 4.8, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 0.72, 1.62, 5.55, 0.34, "Критерий 5 — доступность эталона",
             size=14, bold=True, color=MID)
    text_box(slide, 0.72, 1.96, 5.55, 0.3, "Три уровня причинности Pearl:",
             size=11.5, italic=True, color=DARK_GREY)
    pearl = [
        ("Уровень 3 — контрфактуальность", "«что было бы, если бы X не произошло» — для AI структурно недоступен (территория человека)", SLATE, 0.6),
        ("Уровень 2 — вмешательство", "«что будет, если сделать X» — умеренно (нужен RCT / A/B-тест)", LIGHT, 0.85),
        ("Уровень 1 — ассоциация", "«что бывает вместе с X» — применимость высокая (ML учится на корреляциях)", MID, 1.0),
    ]
    py0 = 2.42; ph = 0.92
    for i, (head, body, col, frac) in enumerate(pearl):
        y = py0 + i * (ph + 0.1)
        w = 5.55 * frac
        rounded_box(slide, 0.72, y, w, ph, fill=WHITE, stroke=col, stroke_w=1.5)
        multiline_box(slide, 0.9, y + 0.08, w - 0.36, ph - 0.16, [
            (head, {"size": 12, "bold": True, "color": col}),
            (body, {"size": 10, "color": DARK_GREY}),
        ], line_spacing=1.04)
    text_box(slide, 0.72, 5.74, 5.55, 0.5,
             "Скорость эталона: компилятор (секунды) → дефолт по кредиту (месяцы) → исход лечения (годы)",
             size=10.5, italic=True, color=MID, line_spacing=1.05)
    # RIGHT: criterion 6 — explainability
    rounded_box(slide, 6.83, 1.5, 6.0, 4.8, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 7.05, 1.62, 5.55, 0.34, "Критерий 6 — объяснимость и аудит",
             size=14, bold=True, color=MID)
    rounded_box(slide, 7.05, 2.05, 5.55, 0.95, fill=BLUE_TINT, stroke=MID, stroke_w=1.25)
    multiline_box(slide, 7.25, 2.16, 5.15, 0.75, [
        ([("EU AI Act ", {"size": 12.5, "bold": True, "color": MID}),
          ("(в силе с 01.08.2024): системы высокого риска обязаны давать объяснимые решения.",
           {"size": 12, "color": DARK_GREY})], {}),
    ], line_spacing=1.08, anchor=MSO_ANCHOR.MIDDLE)
    rounded_box(slide, 7.05, 3.12, 5.55, 2.1, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    multiline_box(slide, 7.25, 3.24, 5.15, 1.85, [
        ("Apple Card 2019", {"size": 12.5, "bold": True, "color": DEEP}),
        ("Вирусная жалоба на 20× разницу лимита; расследование DFS Нью-Йорк закрыто в марте 2021 — Goldman Sachs оправдан в намеренной дискриминации (пол не был во входных признаках).",
         {"size": 11, "color": DARK_GREY}),
    ], line_spacing=1.08)
    gold_callout(slide, 7.05, 5.32, 5.55, 0.88,
                 "Урок — объяснимость, не предвзятость.", size=14)
    add_notes(slide, "s07")
    return slide


def s08(p):
    """assertion_visual — criterion 7 baseline + tools-per-quadrant + UPS photo."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "AI обязан бить лучшую классическую альтернативу на деньгах — иначе отказ "
        "от AI это инженерное решение, не поражение.", y=0.4, size=19)
    # criterion 7 classical examples (left)
    rounded_box(slide, 0.5, 1.5, 7.7, 2.95, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 0.72, 1.62, 7.3, 0.34,
             "Критерий 7 — окупается ли AI против классической альтернативы",
             size=14, bold=True, color=MID)
    items = [
        ("UPS ORION — логистика", "операционные исследования, $300-400 млн экономии/год, ни одной нейросети. ML давал «несколько %», но стоимость интеграции их съедала."),
        ("MPC против RL — нефтехимия", "предиктивное управление работает десятилетиями; RL впервые показан промышленно (Yokogawa FKDPP, 2022), но не вытеснил MPC."),
        ("Запасы — EOQ (Ford-Harris, 1913)", "формула + страховой запас; ML добавляет ценность только на переменном спросе."),
    ]
    iy = 2.02
    for head, body in items:
        multiline_box(slide, 0.72, iy, 7.3, 0.8, [
            (head, {"size": 12.5, "bold": True, "color": DEEP}),
            (body, {"size": 11, "color": DARK_GREY}),
        ], line_spacing=1.05)
        iy += 0.82
    # UPS photo (right)
    rounded_box(slide, 8.4, 1.5, 4.43, 2.95, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    add_image_aspect(slide, SHOTS / "ups-orion.jpg", 8.52, 1.62, 4.19, 2.5)
    img_attribution(slide, 8.52, 4.12, 4.19, "Грузовик UPS · USDA · Public domain (кейс ORION)")
    # tools-per-quadrant (bottom two cols)
    qy = 4.62; qh = 1.55
    rounded_box(slide, 0.5, qy, 6.0, qh, fill=GREEN_TINT, stroke=TEAL, stroke_w=1.5)
    rounded_box(slide, 0.68, qy + 0.14, 5.64, 0.42, fill=TEAL, stroke=None, radius=0.2)
    text_box(slide, 0.8, qy + 0.16, 5.4, 0.38, "Q1 — AI работает + автономно",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, 0.8, qy + 0.66, 5.4, 0.8,
             "Copilot · Claude Code · Stripe Radar · Symbotic · AlphaFold",
             size=12.5, bold=True, color=DEEP, line_spacing=1.1)
    rounded_box(slide, 6.83, qy, 6.0, qh, fill=SOFT_GREY, stroke=SLATE, stroke_w=1.5)
    rounded_box(slide, 7.01, qy + 0.14, 5.64, 0.42, fill=SLATE, stroke=None, radius=0.2)
    text_box(slide, 7.13, qy + 0.16, 5.4, 0.38, "Q3 — «классика выигрывает»",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 7.13, qy + 0.66, 5.4, 0.8, [
        ([("UPS ORION (Gurobi/CPLEX) · MPC · EOQ — ", {"size": 12.5, "color": DEEP}),
          ("AI не нужен", {"size": 13.5, "bold": True, "color": GOLD})], {}),
    ], line_spacing=1.1)
    footer(slide, "Полная стоимость владения включает данные, мониторинг, переобучение, человека в петле, аудит — часто в 3-5× больше лицензии. Измерьте базовую линию до старта.",
           size=10.5)
    add_notes(slide, "s08")
    return slide


def s09(p):
    """assertion_visual — worked example ЖКХ-вода: 7-row verdict checklist."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Вендор обещает «−30% потерь воды за 18 месяцев» — семь критериев за 20 минут "
        "дают вердикт: не полная автоматизация, а L1 advisory + A/B против EPANET.", y=0.4, size=17)
    # scenario box
    rounded_box(slide, 0.5, 1.55, 5.4, 4.7, fill=BLUE_TINT, stroke=MID, stroke_w=1.5)
    multiline_box(slide, 0.72, 1.72, 4.95, 2.2, [
        ("Сценарий", {"size": 14, "bold": True, "color": MID}),
        ("", {"size": 5}),
        ("Вендор «AquaOptima»: deep learning прогноз расхода + автоматическое управление задвижками.",
         {"size": 12, "color": DARK_GREY}),
        ("", {"size": 5}),
        ("Обещание: −30% потерь, 25 млн ₽/год, окупаемость 18 месяцев.",
         {"size": 12, "bold": True, "color": DEEP}),
    ], line_spacing=1.14)
    gold_callout(slide, 0.72, 4.0, 4.95, 2.05,
        "Вердикт: не входить в 6-месячный пилот полной автоматизации. "
        "2 мес. аудит данных + базовая линия EPANET → 4 мес. A/B (AI advisory + EPANET vs EPANET). "
        "Все управляющие действия — человек. Дельта <10% — остаться на EPANET.",
        size=12, anchor=MSO_ANCHOR.TOP)
    # verdict table (right)
    rows = [
        ("1", "Закрытая/открытая", "⚠ полуоткрытая (погода, события, утечки)", WARN, LIGHT),
        ("2", "Данные", "⚠ нужен аудит (утечки не зафиксированы)", WARN, LIGHT),
        ("3", "Повторяемость", "✓ десятки тысяч точек, млн измерений", OK, TEAL),
        ("4", "Цена ошибки / радиус", "✗ задвижка → район без воды на часы", NO, GOLD),
        ("5", "Эталон", "⚠ общий расход — да; локальные утечки — нет", WARN, LIGHT),
        ("6", "Объяснимость", "⚠ публичные решения о ремонтах", WARN, LIGHT),
        ("7", "Экономика vs базовая", "⚠ «−30%» против чего? A/B с EPANET", WARN, LIGHT),
    ]
    tx = 6.1; ty = 1.55; rh = 0.66; tw = 6.73
    for i, (num, crit, verdict, glyph, col) in enumerate(rows):
        y = ty + i * (rh + 0.04)
        is_block = (glyph == NO)
        rounded_box(slide, tx, y, tw, rh,
                    fill=GOLD_TINT if is_block else SURFACE,
                    stroke=GOLD if is_block else SOFT_GREY,
                    stroke_w=1.75 if is_block else 0.75, radius=0.05)
        icon_badge(slide, tx + 0.1, y + 0.13, rh - 0.26, glyph, fill=col, size=16)
        text_box(slide, tx + 0.62, y + 0.04, 2.35, rh - 0.08, f"{num}. {crit}",
                 size=11.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(slide, tx + 3.0, y + 0.04, tw - 3.15, rh - 0.08, verdict,
                 size=10.5, bold=is_block, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    footer(slide, "Критерий 4 (✗) — блокирующий: одна неправильная задвижка оставит район без воды. Эталон — EPANET (открытый гидравлический симулятор водопроводных сетей).",
           size=10.5)
    add_notes(slide, "s09")
    return slide


# ====================================================================
# SECTION 2 — Лестница автономии (s11-s17)
# ====================================================================

def s11(p):
    """schema_layered — ladder L0-L5, bottom-aligned, growing width upward."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Шесть дискретных ступеней — от «человек делает всё» до недостижимой в 2026 "
        "полной автономии.", y=0.4, size=20)
    # ladder bottom-aligned: L0 bottom (narrow), L5 top (widest, grey)
    levels = [
        ("L0 — Без автоматизации", "AI не задействован. Человек делает всё.", "✋", SOFT_GREY, DARK_GREY, 6.0),
        ("L1 — Advisory (советует)", "Классифицирует, предсказывает, рекомендует. Человек решает всегда.", "💡", GOLD_TINT, DEEP, 7.0),
        ("L2 — Supervised (с подтверждением)", "AI действует, человек подтверждает каждое действие.", "✓", TEAL_TINT, DEEP, 8.0),
        ("L3 — Conditional (узкий ODD)", "Автономно в строго ограниченном домене; вне — передаёт человеку.", "◎", BLUE_TINT, DEEP, 9.0),
        ("L4 — High (широкий ODD)", "Действует в широком диапазоне условий; человек на петле.", "◉", BLUE_TINT, DEEP, 10.0),
        ("L5 — Full (любые условия)", "AI решает везде без человека. В 2026 практически недостижим.", "∞", SOFT_GREY, SLATE, 11.0),
    ]
    # draw from bottom (L0) up to top (L5); right column = legend/note
    bottom = 6.35; rh = 0.66; gap = 0.08
    n = len(levels)
    for i, (head, body, glyph, fill, txtcol, w) in enumerate(levels):
        y = bottom - (i + 1) * rh - i * gap
        x = 0.5
        is_l1 = head.startswith("L1")
        is_l5 = head.startswith("L5")
        stroke = GOLD if is_l1 else (SLATE if is_l5 else LIGHT)
        sw = 2.0 if is_l1 else 1.25
        rounded_box(slide, x, y, w, rh, fill=fill, stroke=stroke, stroke_w=sw)
        icon_badge(slide, x + 0.12, y + 0.11, rh - 0.22, glyph,
                   fill=GOLD if is_l1 else MID, size=16)
        multiline_box(slide, x + 0.76, y + 0.05, w - 0.95, rh - 0.08, [
            (head, {"size": 12.5, "bold": True, "color": txtcol}),
            (body, {"size": 10.5, "color": DARK_GREY if not is_l5 else SLATE}),
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    # right note
    rounded_box(slide, 11.3, 1.55, 1.53, 4.8, fill=SURFACE, stroke=LIGHT, stroke_w=1.25)
    up_arrow(slide, 11.85, 2.0, 0.42, 1.6, fill=LIGHT)
    text_box(slide, 11.4, 3.7, 1.33, 2.5,
             "Рост автономии",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.05)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.62,
        "Высокий уровень — не «современнее». Большинство зрелого промышленного AI 2026 года — L1-L2 намеренно.",
        size=12.5)
    add_notes(slide, "s11")
    return slide


def s12(p):
    """schema_matrix — mapping table (3 autonomy scales) + orthogonal axes inset."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Все локальные шкалы курса сводятся в единую L0→L5 — но не все они про "
        "автономию: ось среды и функции ортогональны.", y=0.4, size=18)
    # Table A
    text_box(slide, 0.5, 1.4, 8.0, 0.3, "Таблица A — прямой маппинг (это шкалы автономии)",
             size=13, bold=True, color=MID)
    headers = ["Единая L0→L5", "Разработка ПО", "Авиакосмос", "Производство"]
    cols_x = [0.5, 2.85, 5.0, 7.0]; cols_w = [2.35, 2.15, 2.0, 2.0]
    rows = [
        ("L0 без автомат.", "—", "—", "A0 наблюдать"),
        ("L1 advisory", "A автодополнение", "L1 assistive", "A1 советовать"),
        ("L2 supervised", "B внутри задачи", "L2 supervised", "A2 (часть)"),
        ("L3 conditional", "C код по спецификации", "L3 conditional", "A2 замыкать"),
        ("L4 high", "D инженер-агент", "L4 high", "A3 пилот"),
        ("L5 full", "(недоступен)", "L5 двойн. назнач.", "(редко)"),
    ]
    ty = 1.74
    for j, hdr in enumerate(headers):
        rectangle(slide, cols_x[j], ty, cols_w[j] - 0.05, 0.42, fill=MID)
        text_box(slide, cols_x[j] + 0.06, ty + 0.02, cols_w[j] - 0.16, 0.38, hdr,
                 size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    ry = ty + 0.46; rh = 0.52
    for ri, row in enumerate(rows):
        for j, v in enumerate(row):
            first = (j == 0)
            rounded_box(slide, cols_x[j], ry, cols_w[j] - 0.05, rh,
                        fill=BLUE_TINT if first else SURFACE, stroke=SOFT_GREY,
                        stroke_w=0.75, radius=0.04)
            text_box(slide, cols_x[j] + 0.08, ry + 0.03, cols_w[j] - 0.2, rh - 0.06, v,
                     size=10.5, bold=first, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        ry += rh + 0.04
    # Inset B — orthogonal axes (right column)
    rounded_box(slide, 9.2, 1.4, 3.63, 4.9, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    multiline_box(slide, 9.42, 1.55, 3.2, 4.6, [
        ("Врезка B — ортогональные оси", {"size": 13, "bold": True, "color": DEEP}),
        ("(НЕ автономия)", {"size": 11, "italic": True, "color": DARK_GREY}),
        ("", {"size": 8}),
        ("Логистика — 5 уровней структурированности среды", {"size": 12, "bold": True, "color": MID}),
        ("(контролируемая → город → последняя миля → чёрный лебедь): это ось среды. Чем структурнее среда, тем легче L3-L4.",
         {"size": 10.5, "color": DARK_GREY}),
        ("", {"size": 8}),
        ("Кибербезопасность — «Видит → Решает → Действует»", {"size": 12, "bold": True, "color": MID}),
        ("функциональная декомпозиция: каждая функция может быть на разном уровне автономии.",
         {"size": 10.5, "color": DARK_GREY}),
    ], line_spacing=1.1)
    footer(slide, "Маппинг не точный — это переводчик. Когда вендор говорит «наш AI на L3», спросите: в какой нотации? Источник: SAE J3016 (адаптация).",
           size=10.5)
    add_notes(slide, "s12")
    return slide


def s13(p):
    """assertion_visual — L1 advisory flow + examples + aidoc photo."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "L1 advisory — самый частый уровень зрелой промышленной эксплуатации: "
        "AI советует, человек решает всегда.", y=0.4, size=20)
    # flow: вход → AI артефакт → человек решает
    fy = 1.55; fh = 0.9
    steps = [("Вход", MID), ("AI создаёт\nинформ. артефакт", MID), ("Человек\nрассматривает", LIGHT), ("Человек\nрешает", GOLD)]
    sx = 0.5; sw = 2.55; gap = 0.55
    for i, (label, col) in enumerate(steps):
        x = sx + i * (sw + gap)
        is_last = (i == len(steps) - 1)
        rounded_box(slide, x, fy, sw, fh, fill=GOLD_TINT if is_last else SURFACE,
                    stroke=col, stroke_w=1.75 if is_last else 1.5)
        text_box(slide, x + 0.1, fy + 0.05, sw - 0.2, fh - 0.1, label.replace("\\n", "\n"),
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        if not is_last:
            right_arrow(slide, x + sw + 0.06, fy + fh/2 - 0.16, gap - 0.12, 0.32, fill=LIGHT)
    # examples (left)
    rounded_box(slide, 0.5, 2.7, 7.7, 3.55, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 0.72, 2.82, 7.3, 0.32, "Примеры из курса", size=14, bold=True, color=MID)
    ex = [
        ("Stripe Radar — платежи", "скоринг + флаг «возможный фрод», аналитик решает"),
        ("Aidoc, Chester AI — медицина", "пометка «возможная патология», радиолог решает"),
        ("GitHub Copilot — разработка ПО", "предлагает строку, разработчик принимает"),
        ("Crop Wizard — агро", "отвечает фермеру со ссылками, фермер действует"),
        ("Project Maven — авиакосмос/оборона", "выделяет объекты на снимке, аналитик авторизует"),
    ]
    iy = 3.2
    for head, body in ex:
        multiline_box(slide, 0.72, iy, 7.3, 0.56, [
            ([(head + ": ", {"size": 11.5, "bold": True, "color": DEEP}),
              (body, {"size": 11.5, "color": DARK_GREY})], {}),
        ], line_spacing=1.0)
        iy += 0.58
    # aidoc photo (right top)
    rounded_box(slide, 8.4, 2.7, 4.43, 1.95, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    add_image_aspect(slide, SHOTS / "aidoc.jpg", 8.52, 2.8, 4.19, 1.5)
    img_attribution(slide, 8.52, 4.3, 4.19, "КТ-снимок мозга · Wikimedia · CC BY 4.0 (кейс Aidoc)")
    # criteria up to L2 (right bottom)
    rounded_box(slide, 8.4, 4.75, 4.43, 1.5, fill=BLUE_TINT, stroke=MID, stroke_w=1.5)
    multiline_box(slide, 8.6, 4.86, 4.05, 1.3, [
        ("Подъём на L2 требует:", {"size": 12, "bold": True, "color": MID}),
        ("базовая линия измерена + улучшение AI измерено · контроль изменений · готов откат · приемлемая частота ложных срабатываний.",
         {"size": 10.5, "color": DARK_GREY}),
    ], line_spacing=1.08)
    footer(slide, "Человек решает — всегда. L1 не «недоразвитость», а намеренный режим в регулируемых отраслях.")
    add_notes(slide, "s13")
    return slide


def s14(p):
    """assertion_visual — L2 supervised + L3 conditional, two panels + waymo photo."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "L2 supervised и L3 conditional — AI действует, но человек подтверждает "
        "каждое (L2) или AI работает только в узком домене (L3).", y=0.4, size=19)
    # L2 panel
    rounded_box(slide, 0.5, 1.5, 5.95, 3.4, fill=SURFACE, stroke=TEAL, stroke_w=1.5)
    rounded_box(slide, 0.68, 1.64, 5.59, 0.46, fill=TEAL, stroke=None, radius=0.2)
    text_box(slide, 0.8, 1.66, 5.4, 0.42, "L2 — Supervised (действует с подтверждением)",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.8, 2.22, 5.4, 2.55, [
        ("AI выполняет действие, человек подтверждает или отменяет каждое.",
         {"size": 12, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("Примеры:", {"size": 11.5, "bold": True, "color": MID}),
        ("• Stripe Radar — авто-блок высоко-уверенного фрода (платежи)", {"size": 11, "color": DEEP}),
        ("• Microsoft Defender XDR — авто-карантин (кибербезопасность)", {"size": 11, "color": DEEP}),
        ("• Copilot Workspace — разработка ПО", {"size": 11, "color": DEEP}),
        ("• Yokogawa FKDPP — RL для реактора, оператор вмешивается (процессное производство)", {"size": 11, "color": DEEP}),
    ], line_spacing=1.1)
    # L3 panel
    rounded_box(slide, 6.65, 1.5, 5.95, 3.4, fill=SURFACE, stroke=MID, stroke_w=1.5)
    rounded_box(slide, 6.83, 1.64, 5.59, 0.46, fill=MID, stroke=None, radius=0.2)
    text_box(slide, 6.95, 1.66, 5.4, 0.42, "L3 — Conditional (узкий ODD)",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 6.95, 2.22, 5.4, 2.55, [
        ("Автономно в строго ограниченном домене (ODD); вне домена — отказ и передача человеку. ODD проверяемо в реальном времени.",
         {"size": 12, "color": DARK_GREY}),
        ("", {"size": 4}),
        ("Примеры:", {"size": 11.5, "bold": True, "color": MID}),
        ("• Waymo в геозон-районах (транспорт, фото)", {"size": 11, "color": DEEP}),
        ("• Mobileye Chauffeur на конкретных шоссе (транспорт)", {"size": 11, "color": DEEP}),
        ("• Copilot-агент / Claude Code на очерченном PR (разработка ПО)", {"size": 11, "color": DEEP}),
        ("• See & Spray в узком ODD (агро)", {"size": 11, "color": DEEP}),
    ], line_spacing=1.08)
    # waymo photo (bottom-left)
    rounded_box(slide, 0.5, 5.0, 4.1, 1.55, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    add_image_aspect(slide, SHOTS / "waymo.jpg", 0.62, 5.1, 3.86, 1.18)
    img_attribution(slide, 0.62, 6.25, 3.86, "Waymo robotaxi, SF · Dllu · Wikimedia · CC BY-SA 3.0")
    # ODD + criteria callout (bottom-right)
    gold_callout(slide, 4.85, 5.0, 7.75, 1.55,
        "ODD (Operational Design Domain) — формальное определение условий, где система работает. "
        "Подъём на L4: проверка sim-to-real · канареечный выпуск 5-10% на 2-4 недели · числовые ворота go/no-go.",
        size=12.5)
    add_notes(slide, "s14")
    return slide


def s15(p):
    """assertion_visual — L4 high + L5 blocked, + symbotic photo."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "L4 high достижим в широком домене с человеком на петле; L5 full в 2026 — "
        "теоретический горизонт, не цель.", y=0.4, size=20)
    # L4 (left)
    rounded_box(slide, 0.5, 1.5, 6.0, 4.75, fill=SURFACE, stroke=TEAL, stroke_w=1.5)
    rounded_box(slide, 0.68, 1.64, 5.64, 0.46, fill=TEAL, stroke=None, radius=0.2)
    text_box(slide, 0.8, 1.66, 5.4, 0.42, "L4 — High (широкий ODD, человек на петле)",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    l4 = [
        ("Waymo — транспорт", "сотни тысяч поездок/год без водителя безопасности"),
        ("Symbotic, Amazon Sparrow — логистика (фото)", "миллионы операций в сутки"),
        ("See & Spray в коммерции — агро", "5 млн акров за 3 года"),
        ("AlphaFold (конвейер) — наука", "200 млн структур без вмешательства человека"),
    ]
    iy = 2.24
    for head, body in l4:
        multiline_box(slide, 0.8, iy, 5.4, 0.66, [
            (head, {"size": 12, "bold": True, "color": DEEP}),
            (body, {"size": 11, "color": DARK_GREY}),
        ], line_spacing=1.04)
        iy += 0.68
    # symbotic photo inside L4 panel bottom
    rounded_box(slide, 0.8, 5.0, 5.4, 1.1, fill=WHITE, stroke=LIGHT, stroke_w=1.0)
    add_image_aspect(slide, SHOTS / "symbotic.jpg", 0.9, 5.07, 5.2, 0.78)
    img_attribution(slide, 0.9, 5.82, 5.2, "Amazon fulfilment robot · Wikimedia · CC BY-SA 4.0 (прокси Symbotic)")
    # L5 (right)
    rounded_box(slide, 6.7, 1.5, 6.13, 4.75, fill=SOFT_GREY, stroke=SLATE, stroke_w=1.5)
    rounded_box(slide, 6.88, 1.64, 5.77, 0.46, fill=SLATE, stroke=None, radius=0.2)
    text_box(slide, 7.0, 1.66, 5.5, 0.42, "L5 — Full: почему недостижим в 2026",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    blocks = [
        "1. Страховые рынки отказываются страховать за пределами узких сценариев",
        "2. Регуляторы требуют человеческого надзора для систем высокого риска",
        "3. События вне распределения структурно не лечатся данными",
        "4. Распределение ответственности при L5 юридически не определено",
        "5. Экономика: L4 уже покрывает 99%+ реалистичных сценариев",
    ]
    by = 2.24
    for b in blocks:
        text_box(slide, 7.0, by, 5.5, 0.52, b, size=12, color=DEEP, line_spacing=1.04)
        by += 0.56
    gold_callout(slide, 7.0, 5.12, 5.5, 1.0,
        "L5 — теоретический горизонт, не цель. Tesla позиционируется как «full self-driving», но классификация NHTSA — L2.",
        size=12)
    add_notes(slide, "s15")
    return slide


def s16(p):
    """card_grid 6 — antipatterns per level (2×3 left) + klarna photo (right col)."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "У каждой ступени свой характерный провал — от превышения роли Klarna на L1 "
        "до блока по LAWS на L5.", y=0.4, size=20)
    cards = [
        ("L1 — Превышение роли", "Klarna, поддержка клиентов", "Заявлен советник, действует без подтверждения → обратный найм операторов.", "Позиционирование ≠ режим.", True),
        ("L2 — Скучный человек в петле", "Uber Tempe 2018, транспорт", "Водитель на телефоне, погиб пешеход.", "Скучный мониторинг → HOOL с алертом.", False),
        ("L3 — Расширение ODD без валидации", "Cruise 02.10.2023, транспорт", "Расширили домен без проверки, протащили пешехода 20 футов.", "Расширение ODD требует валидации.", False),
        ("L4 — Действие без канарейки", "CrowdStrike 19.07.2024, кибербез", "Широкий домен = большой радиус, без канареечного выпуска.", "Широкий ODD требует строже.", False),
        ("L5 — Этический / регуляторный блок", "Дебаты о LAWS, авиакосмос/оборона", "Технически возможно, этически и регуляторно — нет.", "L5 для оружия — запрещён.", False),
        ("Сквозной — Пропуск ступени", "Toyota Digit / Cassie, двойники", "Нельзя с L1 на L3 без L2; двойник как мост.", "Каждая ступень — другая дисциплина.", False),
    ]
    cols = 2; cw = 4.52; ch = 1.55; gx = 0.5; gy = 1.55; gapx = 0.2; gapy = 0.13
    for i, (lvl, case, body, lesson, is_l1) in enumerate(cards):
        r, c = divmod(i, cols)
        x = gx + c * (cw + gapx); y = gy + r * (ch + gapy)
        rounded_box(slide, x, y, cw, ch, fill=SURFACE, stroke=GOLD if is_l1 else TEAL, stroke_w=1.75 if is_l1 else 1.5)
        icon_badge(slide, x + 0.14, y + 0.12, 0.36, WARN, fill=GOLD if is_l1 else TEAL, size=13)
        text_box(slide, x + 0.58, y + 0.1, cw - 0.72, 0.4, lvl, size=11, bold=True, color=DEEP, line_spacing=1.0)
        text_box(slide, x + 0.16, y + 0.5, cw - 0.32, 0.26, case, size=9.5, bold=True, italic=True, color=MID)
        text_box(slide, x + 0.16, y + 0.76, cw - 0.32, 0.42, body, size=9.5, color=DARK_GREY, line_spacing=1.0)
        rounded_box(slide, x + 0.14, y + 1.16, cw - 0.28, 0.34, fill=GOLD_TINT, stroke=None, radius=0.14)
        text_box(slide, x + 0.26, y + 1.18, cw - 0.46, 0.3, "Урок: " + lesson, size=9, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # klarna photo column (right)
    px = gx + 2 * cw + gapx + 0.25; pw = 13.333 - px - 0.5
    ph_total = 3 * ch + 2 * gapy
    rounded_box(slide, px, gy, pw, ph_total, fill=WHITE, stroke=GOLD, stroke_w=1.75)
    add_image_aspect(slide, SHOTS / "klarna.jpg", px + 0.12, gy + 0.16, pw - 0.24, ph_total - 0.7)
    img_attribution(slide, px + 0.12, gy + ph_total - 0.5, pw - 0.24,
                    "Стенд Klarna, Internet World Fair 2017 · Wikimedia · CC0 (кейс L1 превышение роли)")
    add_notes(slide, "s16")
    return slide


def s17(p):
    """assertion_visual — worked example exam: ceiling = L1."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "AI-помощник для приёма экзаменов — максимум L1 advisory: высокая цена ошибки "
        "+ обязательная апелляция держат человека в решении.", y=0.4, size=19)
    # compact ladder (left) with ceiling marker on L1
    rounded_box(slide, 0.5, 1.55, 3.7, 4.7, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 0.7, 1.66, 3.3, 0.3, "Лестница — потолок = L1", size=13, bold=True, color=MID)
    lvls = [("L5", SLATE), ("L4", SLATE), ("L3", SLATE), ("L2", SLATE), ("L1", GOLD), ("L0", DARK_GREY)]
    ly = 2.05; lh = 0.62
    for i, (lv, col) in enumerate(lvls):
        y = ly + i * (lh + 0.04)
        is_ceiling = (lv == "L1")
        rounded_box(slide, 0.7, y, 3.3, lh, fill=GOLD_TINT if is_ceiling else WHITE,
                    stroke=GOLD if is_ceiling else SOFT_GREY, stroke_w=1.75 if is_ceiling else 1.0)
        text_box(slide, 0.85, y + 0.03, 1.0, lh - 0.06, lv, size=14, bold=True,
                 color=col if not is_ceiling else DEEP, anchor=MSO_ANCHOR.MIDDLE)
        if is_ceiling:
            text_box(slide, 1.7, y + 0.03, 2.2, lh - 0.06, "← потолок (максимум)", size=11,
                     bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
        elif lv != "L0":
            text_box(slide, 1.7, y + 0.03, 2.2, lh - 0.06, "недопустим", size=10.5,
                     italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
    # analysis (middle)
    rounded_box(slide, 4.4, 1.55, 4.1, 4.7, fill=BLUE_TINT, stroke=MID, stroke_w=1.5)
    multiline_box(slide, 4.62, 1.7, 3.66, 4.4, [
        ("Анализ", {"size": 13, "bold": True, "color": MID}),
        ("", {"size": 6}),
        ("Цена ошибки — высокая (несправедливая оценка, обвинение в плагиате, академический рекорд).",
         {"size": 12, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("Объяснимость — критическая (студент имеет право апеллировать).", {"size": 12, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("Регуляторное окно — академические правила требуют человека-оценщика.", {"size": 12, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("Эталон — для тестов есть, для эссе частично (rubrics, субъективность остаётся).", {"size": 12, "color": DARK_GREY}),
    ], line_spacing=1.08)
    # verdict (right)
    rounded_box(slide, 8.7, 1.55, 4.13, 4.7, fill=SURFACE, stroke=GOLD, stroke_w=1.75)
    gold_callout(slide, 8.9, 1.72, 3.73, 1.5,
        "Вердикт: L1 advisory — максимум. AI выставляет предварительную оценку + флаг плагиата + объяснение; преподаватель решает; студент может оспорить.",
        size=12, anchor=MSO_ANCHOR.TOP)
    multiline_box(slide, 8.9, 3.4, 3.73, 2.7, [
        ("Подъём на L2:", {"size": 12.5, "bold": True, "color": MID}),
        ("измерить базовую линию (сколько ошибается преподаватель / AI на тест-сете); если AI систематически точнее — переход с контролем изменений.",
         {"size": 11.5, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("L3-L4 недопустимы по регулятору и этике.", {"size": 11.5, "bold": True, "color": DEEP}),
    ], line_spacing=1.1)
    add_notes(slide, "s17")
    return slide


# ====================================================================
# SECTION 3 — Карта 16 отраслей (s19-s26)
# ====================================================================

def s19(p):
    """lecture_map — 4×4 grid, industry + local axis (no L-codes)."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Каждая из 16 отраслей курса поднимала свою локальную ось — и все они про "
        "одно: насколько AI берёт действие на себя.", y=0.4, size=19)
    cards = [
        ("Введение в AI", "Диагностический вопрос; эффект ИИ; 3 уровня Pearl", MID),
        ("Большие модели", "Конвейер вывода; «слепота к буквам» как граница", MID),
        ("Архитектуры AI", "prompt → RAG → агент → дообучение; накопление ненадёжности", MID),
        ("Разработка ПО", "Лестница A/B/C/D (автодополнение → инженер-агент)", MID),
        ("Финансы и ритейл", "Закрытый мир vs открытый мир прогноза", MID),
        ("CAD/CAM", "6 классов AI в проектировании; «оптимизация ≠ генерация»", MID),
        ("Медицина", "Закрытая петля; человек в петле; RCT как эталон", MID),
        ("Креатив", "«Добавил → изменил → сломал» × 4 области", MID),
        ("Авиакосмос/оборона", "OODA; человек в/на/вне петли; автономия L1-L5", TEAL),
        ("Агросектор", "Лестница АПК; закрытая петля vs открытая среда", TEAL),
        ("Производство", "Дискретное vs процессное; болото пилотов", TEAL),
        ("Цифровые двойники", "A0-A3 + цифровой двойник как мост", TEAL),
        ("Логистика/транспорт", "Лестница среды (5 уровней); ODD; 7 критериев", GOLD),
        ("Телеком/кибербез", "«Видит → Решает → Действует»; MITRE ATLAS", GOLD),
        ("Наука", "Лестница научного цикла; AlphaFold vs Galactica", GOLD),
        ("Нефтегаз", "Матрица 2×2 «данные × процесс»", GOLD),
    ]
    cols, rows = 4, 4
    gx, gy = 0.5, 1.5
    cw, ch = 3.04, 1.2
    pad = 0.07
    for i, (name, axis, col) in enumerate(cards):
        r, c = divmod(i, cols)
        x = gx + c * (cw + pad); y = gy + r * (ch + pad)
        rounded_box(slide, x, y, cw, ch, fill=SURFACE, stroke=col, stroke_w=1.75)
        text_box(slide, x + 0.14, y + 0.08, cw - 0.28, 0.36, name, size=13, bold=True, color=col, line_spacing=1.0)
        text_box(slide, x + 0.14, y + 0.44, cw - 0.28, ch - 0.5, axis, size=10, color=DARK_GREY, line_spacing=1.04)
    footer(slide, "16 локальных осей автономии / применимости — переименования одного явления, а не разные явления. Capstone-карта объединяет их все.")
    add_notes(slide, "s19")
    return slide


def scatter_slide(p, assertion, chart_name, side_blocks, notes_key, *,
                  asize=20, footer_text=None, gold_block=None):
    """Generic scatter slide: chart left (~64%), annotation column right."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide, assertion, y=0.4, size=asize)
    rounded_box(slide, 0.5, 1.4, 8.5, 5.0, fill=WHITE, stroke=LIGHT, stroke_w=1.5)
    add_image_aspect(slide, CHARTS / chart_name, 0.62, 1.52, 8.26, 4.76)
    # right annotation column
    rx = 9.2; rw = 3.63
    rounded_box(slide, rx, 1.4, rw, 5.0, fill=SURFACE, stroke=LIGHT)
    iy = 1.6
    for head, body, col in side_blocks:
        multiline_box(slide, rx + 0.2, iy, rw - 0.4, 0.95, [
            (head, {"size": 12.5, "bold": True, "color": col}),
            (body, {"size": 10.5, "color": DARK_GREY}),
        ], line_spacing=1.06)
        iy += 1.02
    if gold_block:
        gold_callout(slide, rx + 0.2, 5.45, rw - 0.4, 0.85, gold_block, size=11.5,
                     anchor=MSO_ANCHOR.MIDDLE)
    if footer_text:
        footer(slide, footer_text, size=10.5)
    add_notes(slide, notes_key)
    return slide


def s20(p):
    return scatter_slide(p,
        "Первые четыре точки уже показывают структуру: смежное с IT — вверху-справа, "
        "регулируемая медицина — внизу-справа.", "s20-batch1.png", [
            ("Разработка ПО — верх-право", "текст/код — модальность фундаментальных моделей; компилятор даёт эталон; объём гигантский.", MID),
            ("Финансы и ритейл — верх-середина", "фрод (закрытый мир) высоко; iBuying (открытый мир) низко — Zillow показал почему.", MID),
            ("Медицина — низ-право", "узкая визуализация высоко; автономия капнута FDA на L1.", MID),
            ("Авиакосмос/оборона — потолок", "восприятие высоко, действие ограничено этически (дебаты о LAWS).", TEAL),
        ], "s20", asize=20,
        gold_block="Разработка ПО — чистая координата верх-право.")


def s21(p):
    return scatter_slide(p,
        "Средний набор добавляет двойственные (бимодальные) отрасли: агро раздваивается "
        "на узкую очерченную See & Spray и провальный Monarch.", "s21-batch2.png", [
            ("CAD/CAM — середина", "оптимизация высоко; LLM для CAD-скриптов низко (ORCA 45-63%).", MID),
            ("Креатив — бимодально", "массовые ассеты высоко; авторское произведение низко (Getty vs Stability).", MID),
            ("Агро — двойственное облако", "See & Spray ↑ верх-право; Monarch ↖ верх-лево (зона предупреждения), провал.", TEAL),
            ("Производство — середина-низ", "CV на конвейере высоко; LLM на коде PLC низко.", TEAL),
        ], "s21", asize=18,
        gold_block="Агро раздваивается — готовность локальна по задачам, не глобальна по отрасли.")


def s22(p):
    return scatter_slide(p,
        "Логистика — три точки одной отрасли: склад L4, робот-такси L3, чёрный лебедь "
        "фактически L0.", "s22-batch3-full.png", [
            ("Логистика (3 точки!)", "склад (Symbotic, L4) ↑; робот-такси (Waymo/Cruise) середина; чёрный лебедь (Суэц, COVID) — фактически L0.", GOLD),
            ("Телеком/кибербез", "восприятие высоко (L1-L3); действие ограничено радиусом (CrowdStrike).", TEAL),
            ("Наука — двойственная", "AlphaFold (закрытый мир) высоко vs Galactica (открытый мир) низко.", GOLD),
            ("Нефтегаз", "LLM без физики не заменит геолога — нужны PINN / гибрид физика+ML.", GOLD),
        ], "s22", asize=20,
        footer_text="Логистика — самая бимодальная: внутри одной компании AI на L4 (склад) и на L0 (чёрный лебедь). Зрелость локальна, не глобальна.")


def cluster_slide(p, assertion, chart_name, traits_title, traits, lesson, notes_key,
                  *, photo=None, photo_attr=None, asize=20, gold_lesson=True):
    """Cluster analysis slide: highlighted-quadrant scatter left, traits right."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide, assertion, y=0.4, size=asize)
    # scatter left
    rounded_box(slide, 0.5, 1.45, 6.3, 4.95, fill=WHITE, stroke=LIGHT, stroke_w=1.5)
    add_image_aspect(slide, CHARTS / chart_name, 0.6, 1.55, 6.1, 4.75)
    # traits right
    rx = 7.0; rw = 5.83
    rounded_box(slide, rx, 1.45, rw, 3.55, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, rx + 0.22, 1.58, rw - 0.44, 0.34, traits_title, size=14, bold=True, color=MID)
    ty = 1.98
    for t in traits:
        text_box(slide, rx + 0.22, ty, rw - 0.44, 0.5, t, size=11.5, color=DEEP, line_spacing=1.04)
        ty += 0.5
    # lesson / photo bottom
    if photo:
        rounded_box(slide, rx, 5.1, 2.0, 1.3, fill=WHITE, stroke=LIGHT, stroke_w=1.0)
        add_image_aspect(slide, SHOTS / photo, rx + 0.08, 5.16, 1.84, 0.92)
        img_attribution(slide, rx + 0.08, 6.1, 1.84, photo_attr)
        lx = rx + 2.15; lw = rw - 2.15
    else:
        lx = rx; lw = rw
    if gold_lesson:
        gold_callout(slide, lx, 5.1, lw, 1.3, lesson, size=12, anchor=MSO_ANCHOR.MIDDLE)
    else:
        rounded_box(slide, lx, 5.1, lw, 1.3, fill=BLUE_TINT, stroke=MID, stroke_w=1.5)
        text_box(slide, lx + 0.2, 5.18, lw - 0.4, 1.14, lesson, size=12, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    add_notes(slide, notes_key)
    return slide


def s23(p):
    return cluster_slide(p,
        "Кластер закрытой петли (верхний правый) — пять общих признаков, профиль "
        "успешного применения AI.", "s23-cluster-ur.png",
        "Пять общих признаков — профиль успеха", [
            "1. Среда контролируемая или закрытая (критерий 1 ✓)",
            "2. Эталон быстрый и однозначный (критерий 5 ✓): компилятор / чарджбэк / разметка",
            "3. Объём задач большой (критерий 3 ✓)",
            "4. Цена ошибки низкая или поглощаемая (критерий 4 ✓)",
            "5. Распределение обучения совпадает с эксплуатацией (критерий 2 ✓)",
        ],
        "Новая отрасль удовлетворяет всем пяти признакам — у проекта хорошие шансы.",
        "s23", photo="alphafold.png",
        photo_attr="AlphaFold · Wikimedia · CC BY-SA 4.0", asize=20)


def s24(p):
    return cluster_slide(p,
        "Кластер открытой среды (верхний левый) — зона провалов; открытая среда это "
        "физика задачи, а не недостаток модели.", "s24-cluster-ul.png",
        "Состав и что общего (зона предупреждения)", [
            "Monarch (38% сокращений ≈53/140) · Plenty ($940 млн+ → Chapter 11)",
            "Cruise · чёрный лебедь (Suez, COVID) · Zillow ($304 млн) · Galactica (отозвана за 48 ч)",
            "1. Среда открытая, состязательная (критерий 1 ✗)",
            "2. Данные не покрывают редкие события (критерий 2 ✗)",
            "3. Цена ошибки высокая (критерий 4 ✗) · 4. Эталон медленный/нет (критерий 5 ✗)",
        ],
        "Открытая среда — это физика задачи, а не недостаток модели. Сузьте домен; перейдите на advisory; откажитесь от AI.",
        "s24", photo="monarch.jpg",
        photo_attr="Monarch MK-V · monarchtractor.com · fair use", asize=19)


def s25(p):
    return cluster_slide(p,
        "Кластер высоких ставок (нижний правый) — применимость высокая, но автономия "
        "капнута регулятором; цель — усиление, не замена.", "s25-cluster-lr.png",
        "Состав (регуляторно-капнутая автономия)", [
            "Медицина (капнут FDA на L1) · Авиакосмос (капнут регуляцией LAWS)",
            "Производство, критичное к безопасности (IEC 61508 SIL 2/3)",
            "Кибербезопасность уровня действия (дисциплина радиуса)",
            "Нефтегаз ATEX Zone 0 (сертификация оборудования)",
            "Общее: применимость высокая, автономия ограничена; объяснимость обязательна.",
        ],
        "Усиление, а не замена. Радиолог с AI лучше радиолога без — но решение всегда за радиологом (IBM Watson Health показал, что полная замена проваливается).",
        "s25", photo="epic-sepsis.jpg",
        photo_attr="Мониторинг ICU · US Navy · Public domain (кейс Epic Sepsis)", asize=19)


def s26(p):
    """schema_quadrant — empty upper-left + shift arrows."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Пустой верхне-левый квадрант — дидактический сигнал: низкая применимость "
        "× высокая автономия = катастрофа.", y=0.4, size=20)
    rounded_box(slide, 0.5, 1.45, 7.4, 4.95, fill=WHITE, stroke=LIGHT, stroke_w=1.5)
    add_image_aspect(slide, CHARTS / "s26-empty-quadrants.png", 0.6, 1.55, 7.2, 4.75)
    # right column
    rx = 8.1; rw = 4.73
    rounded_box(slide, rx, 1.45, rw, 2.5, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.75)
    multiline_box(slide, rx + 0.22, 1.6, rw - 0.44, 2.25, [
        ("Верхний левый — пустой и опасный", {"size": 13, "bold": True, "color": DEEP}),
        ("Кто туда попадает:", {"size": 11, "italic": True, "color": DARK_GREY}),
        ("• CrowdStrike Falcon BSOD — радиус 8,5 млн устройств", {"size": 11, "color": DEEP}),
        ("• F-35 ALIS — предиктивное обслуживание с высокой автономией", {"size": 11, "color": DEEP}),
        ("• Cruise — расширение домена в открытой среде", {"size": 11, "color": DEEP}),
    ], line_spacing=1.1)
    rounded_box(slide, rx, 4.05, rw, 2.35, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    multiline_box(slide, rx + 0.22, 4.18, rw - 0.44, 2.1, [
        ("Два направления сдвига:", {"size": 12.5, "bold": True, "color": MID}),
        ("", {"size": 4}),
        ([("Вправо → ", {"size": 12, "bold": True, "color": TEAL}),
          ("повысить применимость (закрыть среду, набрать данные, улучшить эталон).", {"size": 11.5, "color": DARK_GREY})], {}),
        ("", {"size": 4}),
        ([("Вниз ↓ ", {"size": 12, "bold": True, "color": TEAL}),
          ("снизить автономию (на-петле → в-петле → советник).", {"size": 11.5, "color": DARK_GREY})], {}),
    ], line_spacing=1.12)
    footer(slide, "Попадаю ли я в верхне-левый? Если да — что изменить, чтобы сдвинуться?")
    add_notes(slide, "s26")
    return slide


# ====================================================================
# SECTION 4 — Двенадцать провалов (s28-s31)
# ====================================================================

def failure_grid(slide, cards, *, photo=None, photo_attr=None, gold_idx=None):
    """2×2 failure card grid (left). cards: (name, sources, lesson, alt).
    Optional real photo + caption in a right column (~3.0")."""
    cols = 2; cw = 4.52; ch = 2.35; gx = 0.5; gy = 1.5; gapx = 0.2; gapy = 0.18
    for i, (name, sources, lesson, alt) in enumerate(cards):
        r, c = divmod(i, cols)
        x = gx + c * (cw + gapx); y = gy + r * (ch + gapy)
        is_gold = (gold_idx == i)
        rounded_box(slide, x, y, cw, ch, fill=SURFACE, stroke=GOLD if is_gold else TEAL,
                    stroke_w=1.85 if is_gold else 1.5)
        icon_badge(slide, x + 0.14, y + 0.14, 0.4, WARN, fill=GOLD if is_gold else TEAL, size=14)
        text_box(slide, x + 0.62, y + 0.13, cw - 0.76, 0.5, name, size=11.5, bold=True, color=DEEP, line_spacing=1.0)
        text_box(slide, x + 0.16, y + 0.66, cw - 0.32, 0.36, sources, size=9.5, italic=True, color=MID, line_spacing=1.0)
        multiline_box(slide, x + 0.16, y + 1.0, cw - 0.32, 0.66, [
            ([("Урок: ", {"size": 10, "bold": True, "color": DEEP}),
              (lesson, {"size": 10, "color": DARK_GREY})], {}),
        ], line_spacing=1.0)
        rounded_box(slide, x + 0.14, y + 1.66, cw - 0.28, 0.6, fill=GREEN_TINT, stroke=None, radius=0.1)
        multiline_box(slide, x + 0.26, y + 1.69, cw - 0.46, 0.54, [
            ([("Альтернатива: ", {"size": 9.5, "bold": True, "color": DEEP}),
              (alt, {"size": 9.5, "bold": True, "color": DEEP})], {}),
        ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # right photo column
    if photo:
        px = gx + 2 * cw + gapx + 0.25; pw = 13.333 - px - 0.5
        ph_total = ch * 2 + gapy
        rounded_box(slide, px, gy, pw, ph_total, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
        add_image_aspect(slide, SHOTS / photo, px + 0.12, gy + 0.16, pw - 0.24, ph_total - 0.62)
        img_attribution(slide, px + 0.12, gy + ph_total - 0.42, pw - 0.24, photo_attr)


def s28(p):
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Провалы 1-4: открытый мир, накопление ненадёжности, демо ≠ промышленная "
        "эксплуатация, скучный человек в петле.", y=0.4, size=19)
    cards = [
        ("1. Прогноз в открытом мире без закрытой петли", "Zillow · Monarch · Cruise",
         "распределённый сдвиг убивает ML на исторических данных.", "сузить домен; советник; не-AI базовая линия."),
        ("2. Накопление ненадёжности в многошаговом агенте", "$4 200-петля у агента поддержки; агентная разработка ПО",
         "p^N → 0 при N>10 (p=0,95, N=10 → 0,59).", "лимит бюджета + макс. шагов + точки человека."),
        ("3. Демо вендора ≠ промышленная эксплуатация", "Devin · IBM Watson · Epic Sepsis · Klarna",
         "вендорский бенчмарк ≠ ваша среда.", "повторить измерение на ваших данных."),
        ("4. Скучный человек в петле не работает", "Uber Tempe · F-35 ALIS",
         "редкий скучный мониторинг проваливается.", "человек-на-петле с алертом."),
    ]
    failure_grid(slide, cards, gold_idx=1, photo="uber-tempe.jpg",
                 photo_attr="Uber Volvo XC90 (модель из аварии Tempe 2018) · Dllu · Wikimedia · CC BY-SA 2.0")
    footer(slide, "Источники кейсов — отраслевые лекции курса. Главное — математика накопления ненадёжности p^N → 0.", size=10.5)
    add_notes(slide, "s28")
    return slide


def s29(p):
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Провалы 5-8: переавтоматизация, действие без канарейки, научная галлюцинация, "
        "голос/видео-дипфейк мошенничество.", y=0.4, size=19)
    cards = [
        ("5. Переавтоматизация в зонах вариативности", "Tesla 2018 · Boeing 737 MAX 9",
         "парадокс автоматизации (Бейнбридж, 1983): автоматизация ломает пропускную способность.", "Jidoka — усиление; SPC остаётся."),
        ("6. Действие без канарейки и отката", "CrowdStrike · Cloudflare",
         "широкий домен = большой радиус разрушения.", "канарейка 1-5% + телеметрия + откат в один клик."),
        ("7. Научная галлюцинация класса Galactica", "Meta Galactica — галлюцинации цитат",
         "в науке эталон — воспроизводимый эксперимент, а не текст.", "RAG-заземление + рецензирование по чек-листу."),
        ("8. Голос/чат-мошенничество и дипфейк", "Wendy's · Air Canada · Arup $25 млн (видео-конф., февр. 2024)",
         "шумная среда + сложность = AI ненадёжен; видео+голос = новый вектор.", "C2PA-подпись; проверка через независимый канал."),
    ]
    failure_grid(slide, cards, gold_idx=3, photo="arup-deepfake.jpg",
                 photo_attr="Видеоконференция (концепт; кейс Arup $25 млн) · DoD · CC BY 2.0")
    footer(slide, "Самая яркая цифра новых векторов мошенничества — дипфейк финансового директора Arup на $25 млн.", size=10.5)
    add_notes(slide, "s29")
    return slide


def s30(p):
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Провалы 9-12: утечка обучающих данных, привязка к вендору, slopsquatting, "
        "болото пилотов.", y=0.4, size=19)
    cards = [
        ("9. Дословная утечка обучающих данных", "Getty vs Stability · NYT vs OpenAI",
         "у фундаментальных моделей есть хвост запоминания — часть данных хранится дословно.", "лицензированные датасеты; аудит провенанса; C2PA."),
        ("10. Привязка к вендору в регулируемых отраслях", "Climate FieldView · F-35 ALIS · IBM Watson",
         "привязка → стратегический риск; данные у вендора, выход стоит миллионы.", "гос-инфраструктура; экспорт данных в контракте."),
        ("11. Slopsquatting / галлюцинации цепочки поставок", "Атаки на имена пакетов npm/pip — разработка ПО",
         "AI выдумывает имена пакетов; атакующие их занимают → вредоносный код.", "проверка SBOM + список разрешённых импортов."),
        ("12. Болото пилотов / 90-95% не доходят", "MIT NANDA 95% · McKinsey 5,5%* · РФ 9 из 10 · двойники 75%",
         "пилот соскальзывает в бесконечность.", "явные ворота GO/NO-GO + базовая линия + лимит бюджета."),
    ]
    failure_grid(slide, cards, gold_idx=3, photo="getty-stability.jpg",
                 photo_attr="Stable Diffusion с искажённым Getty-watermark (кейс Getty v. Stability) · AI-gen · Public domain")
    footer(slide, "Главная статистика курса — 90-95% пилотов не доходят. *95% (MIT) и 5,5% (McKinsey) — РАЗНЫЕ измерения, не одно число.", size=10.5)
    add_notes(slide, "s30")
    return slide


def s31(p):
    """assertion_visual — 3 mega-patterns + 30-sec procedure."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Почти все провалы сводятся к трём мега-паттернам: AI за границей закрытой петли, "
        "плохой человек в петле, проигнорированная базовая линия.", y=0.4, size=19)
    mega = [
        ("Мега-паттерн 1", "AI за границей закрытой петли", "Провалы 1 (Zillow/Monarch/Cruise), 7 (Galactica), частично 5.", "Диагностируйте среду до старта (критерий 1).", "↻"),
        ("Мега-паттерн 2", "Человек в петле спроектирован плохо", "Провалы 4 (Uber Tempe/F-35 ALIS), частично 1 (Cruise), 6 (CrowdStrike).", "Человек в петле — инженерная дисциплина, не формулировка.", "◉"),
        ("Мега-паттерн 3", "Экономическая базовая линия проигнорирована", "Провалы 3 (демо ≠ эксплуатация), 12 (болото пилотов), частично 10.", "Измерьте базовую линию + классическую альтернативу до коммита.", "⚖"),
    ]
    cw = 4.04; gx = 0.5; gy = 1.55; ch = 3.4; gapx = 0.13
    for i, (tag, title, body, lesson, glyph) in enumerate(mega):
        x = gx + i * (cw + gapx)
        rounded_box(slide, x, gy, cw, ch, fill=SURFACE, stroke=MID, stroke_w=1.5)
        icon_badge(slide, x + cw/2 - 0.32, gy + 0.22, 0.64, glyph, fill=MID, size=22)
        text_box(slide, x + 0.16, gy + 0.98, cw - 0.32, 0.3, tag, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        text_box(slide, x + 0.18, gy + 1.3, cw - 0.36, 0.7, title, size=13.5, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.04)
        text_box(slide, x + 0.18, gy + 2.0, cw - 0.36, 0.7, body, size=10.5, color=DARK_GREY, align=PP_ALIGN.CENTER, line_spacing=1.04)
        rounded_box(slide, x + 0.16, gy + 2.72, cw - 0.32, 0.56, fill=GOLD_TINT, stroke=None, radius=0.1)
        text_box(slide, x + 0.28, gy + 2.75, cw - 0.5, 0.5, lesson, size=10, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # 30-sec procedure bar
    rounded_box(slide, 0.5, 5.25, 12.33, 1.05, fill=BLUE_TINT, stroke=MID, stroke_w=1.5)
    text_box(slide, 0.72, 5.34, 3.4, 0.9, "30-секундная процедура\nдля любого AI-предложения:",
             size=13, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    multiline_box(slide, 4.2, 5.36, 8.4, 0.85, [
        ("1. Какая среда — закрытая или открытая со сдвигом?", {"size": 12, "color": DEEP}),
        ("2. Есть ли человек в петле, и не скучно ли ему?", {"size": 12, "color": DEEP}),
        ("3. Какая базовая линия, и окупит ли AI дельту?", {"size": 12, "color": DEEP}),
    ], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "s31")
    return slide


# ====================================================================
# SECTION 5 — Карточки (s32-s36)
# ====================================================================

def s32(p):
    """assertion_visual — overview of 4 cheat-sheets."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Четыре опорные карточки — то, что вы заберёте с собой и будете применять "
        "в работе.", y=0.42, size=22)
    cards = [
        ("#1 (A4)", "Матрица из 7 критериев", "«применять ли AI» — чек-лист", "▤", MID, False),
        ("#2 (A4)", "Лестница автономии L0→L5", "+ критерии подъёма", "▥", TEAL, False),
        ("#3 (A4)", "Реестр 12 провалов", "+ противоядия", "⚠", GOLD, True),
        ("#4 (A1)", "Карта 16 отраслей", "плакат на 2D-плоскости", "◫", GOLD, False),
    ]
    cw = 2.95; gx = 0.5; gy = 2.0; ch = 3.6; gap = 0.18
    for i, (num, title, sub, glyph, col, is_main) in enumerate(cards):
        x = gx + i * (cw + gap)
        rounded_box(slide, x, gy, cw, ch, fill=GOLD_TINT if is_main else SURFACE,
                    stroke=GOLD if is_main else col, stroke_w=2.25 if is_main else 1.5)
        icon_badge(slide, x + cw/2 - 0.42, gy + 0.35, 0.84, glyph, fill=col, size=26)
        text_box(slide, x + 0.12, gy + 1.4, cw - 0.24, 0.4, "Карточка " + num, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        text_box(slide, x + 0.14, gy + 1.85, cw - 0.28, 0.85, title, size=14.5, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.05)
        text_box(slide, x + 0.14, gy + 2.7, cw - 0.28, 0.6, sub, size=11.5, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER, line_spacing=1.05)
        if is_main:
            chip(slide, x + cw/2 - 1.0, gy + ch - 0.5, 2.0, 0.36, "главная по ценности", fill=GOLD, color=DEEP, size=10)
    footer(slide, "Карточки доступны как PDF + исходники в репозитории курса. Распечатайте; повесьте #3 рядом со столом, #4 как большой плакат в лаборатории.")
    add_notes(slide, "s32")
    return slide


def cheatsheet_table(slide, headers, rows, col_w, *, header_y=1.5, rh=0.6,
                     footer_text=None, gold_footer=False, qr=False):
    """Render a fullscreen-ish cheat-sheet table preview."""
    total_w = sum(col_w)
    x0 = (13.333 - total_w) / 2
    # headers
    cx = x0
    for j, hdr in enumerate(headers):
        rectangle(slide, cx, header_y, col_w[j] - 0.05, 0.46, fill=MID)
        text_box(slide, cx + 0.08, header_y + 0.03, col_w[j] - 0.2, 0.4, hdr,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    ry = header_y + 0.5
    for ri, row in enumerate(rows):
        cx = x0
        tint = SURFACE if ri % 2 == 0 else WHITE
        for j, v in enumerate(row):
            rounded_box(slide, cx, ry, col_w[j] - 0.05, rh, fill=tint, stroke=SOFT_GREY,
                        stroke_w=0.6, radius=0.03)
            text_box(slide, cx + 0.08, ry + 0.03, col_w[j] - 0.2, rh - 0.06, v,
                     size=11, bold=(j == 0), color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
            cx += col_w[j]
        ry += rh + 0.03
    if footer_text:
        if gold_footer:
            gold_callout(slide, x0, ry + 0.08, total_w - (1.0 if qr else 0), 0.62, footer_text, size=12)
        else:
            footer(slide, footer_text)
    if qr:
        rounded_box(slide, x0 + total_w - 0.85, ry + 0.08, 0.7, 0.62, fill=WHITE, stroke=LIGHT, stroke_w=1.0)
        text_box(slide, x0 + total_w - 0.85, ry + 0.16, 0.7, 0.5, "QR\nPDF", size=9, color=SLATE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def s33(p):
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide, "Карточка #1 — матрица из 7 критериев: пройди 7 строк, один ✗ — СТОП.",
                    y=0.42, size=21)
    headers = ["№", "Вопрос", "Индикатор", "Пример из курса"]
    col_w = [0.7, 4.0, 3.6, 3.7]
    rows = [
        ("1", "Среда контролируемая/закрытая?", "✓ закрытая / ⚠ полу / ✗ открытая", "See & Spray ✓; Monarch ✗"),
        ("2", "Данных достаточно + совпадают?", "✓ / ⚠ / ✗", "компилятор ✓; Epic Sepsis ✗"),
        ("3", "Повторяема + высокий объём?", "✓ / ✗", "Copilot ✓; штучная архитектура ✗"),
        ("4", "Цена ошибки приемлема?", "✓ низкая / ⚠ HITL / ✗ катастрофа", "Stripe Radar ✓; CrowdStrike ✗"),
        ("5", "Эталон быстрый?", "✓ / ✗", "компилятор ✓; скупка жилья ✗"),
        ("6", "Нужна объяснимость?", "✓ SHAP/LIME / ⚠ прозрачн. модель / ✗", "Aidoc; Apple Card 2019"),
        ("7", "Окупается vs базовая линия?", "✓ / ✗", "LaserWeeder ✓; ML vs MPC ✗"),
    ]
    cheatsheet_table(slide, headers, rows, col_w, header_y=1.55, rh=0.6,
        footer_text="Один ✗ — СТОП, отказ от полного AI. ≥2 ⚠ — СТОП, обоснуй человека в петле + канарейка + откат. Все 7 ✓ — пилот с явными GO/NO-GO.",
        gold_footer=True, qr=True)
    add_notes(slide, "s33")
    return slide


def s34(p):
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide, "Карточка #2 — лестница автономии L0→L5: определи текущий уровень и максимально допустимый.",
                    y=0.42, size=19)
    headers = ["Уровень", "Название", "Что AI делает", "Кто решает", "Критерий подъёма"]
    col_w = [1.05, 1.85, 2.85, 2.4, 3.6]
    rows = [
        ("L0", "Без автоматизации", "Нет AI", "Человек", "Базовая линия собрана"),
        ("L1", "Advisory", "Классиф./предсказ./реком.", "Человек всегда", "База + контроль изменений + откат"),
        ("L2", "Supervised", "Действует, человек подтв.", "Человек каждое", "Ложные срабат. + канарейка + откат"),
        ("L3", "Conditional", "Действует в узком ODD", "На петле (HOOL)", "ODD + телеметрия + ворота go/no-go"),
        ("L4", "High", "Действует в широком ODD", "Вне петли (большинство)", "99,9% + страховка + допуск регулятора"),
        ("L5", "Full", "Решает везде", "(недостижим 2026)", "Для большинства отраслей недоступен"),
    ]
    cheatsheet_table(slide, headers, rows, col_w, header_y=1.5, rh=0.6,
        footer_text="Антипаттерны: L1 превышение роли (Klarna) · L2 скучный HITL (Uber Tempe) · L3 расширение ODD (Cruise) · L4 без канарейки (CrowdStrike) · L5 этический блок (LAWS) · пропуск ступени (двойник как мост).",
        gold_footer=False)
    add_notes(slide, "s34")
    return slide


def s35(p):
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide, "Карточка #3 — реестр 12 провалов: главная по ценности, прививка против реальных ошибок.",
                    y=0.4, size=19)
    headers = ["№", "Провал", "Источник", "Урок → альтернатива"]
    col_w = [0.55, 3.5, 3.1, 5.05]
    rows = [
        ("1", "Открытый мир без петли", "Zillow/Monarch/Cruise", "сдвиг убивает ML → сузить домен"),
        ("2", "Накопление ненадёжности", "$4 200-петля / разработка ПО", "p^N→0 при N>10 → лимит + точки HITL"),
        ("3", "Демо ≠ эксплуатация", "Devin/Watson/Epic/Klarna", "бенчмарк ≠ ваша среда → повторить на своих"),
        ("4", "Скучный HITL", "Uber Tempe / F-35 ALIS", "мониторинг проваливается → HOOL с алертом"),
        ("5", "Переавтоматизация", "Tesla 2018 / Boeing MAX 9", "парадокс Bainbridge → Jidoka, усиление"),
        ("6", "Действие без канарейки", "CrowdStrike / Cloudflare", "широкий ODD = радиус → канарейка + откат"),
        ("7", "Научная галлюцинация", "Galactica", "текст ≠ эксперимент → RAG + рецензирование"),
        ("8", "Голос/видео-дипфейк", "Wendy's / Air Canada / Arup $25M", "шум+сложность → C2PA + проверка"),
        ("9", "Утечка обуч. данных", "Getty / NYT", "хвост запоминания → лицензир. датасеты"),
        ("10", "Привязка к вендору", "Climate FieldView/ALIS/Watson", "стратегич. риск → экспорт данных"),
        ("11", "Slopsquatting", "имена npm / pip", "выдуманные имена = атака → SBOM + allow-list"),
        ("12", "Болото пилотов", "MIT 95% / McKinsey 5,5% / РФ 9 из 10", "слип в бесконечность → ворота GO/NO-GO"),
    ]
    cheatsheet_table(slide, headers, rows, col_w, header_y=1.4, rh=0.41,
        footer_text="Узнаёшь паттерн — задай уточняющий вопрос. Часто этого достаточно, чтобы спасти проект.",
        gold_footer=True, qr=True)
    add_notes(slide, "s35")
    return slide


def s36(p):
    """cheatsheet_preview — A1 master poster (full scatter)."""
    slide = blank(p); set_slide_bg(slide, WHITE)
    assertion_title(slide,
        "Карточка #4 — плакат A1: карта 16 отраслей на стену; для новой отрасли — "
        "аналогии по координатам.", y=0.4, size=19)
    rounded_box(slide, 0.5, 1.4, 12.33, 4.95, fill=WHITE, stroke=LIGHT, stroke_w=1.5)
    add_image_aspect(slide, CHARTS / "s36-master-poster.png", 0.65, 1.52, 12.03, 4.7)
    footer(slide, "Повесьте на стену. Для новой отрасли: «эта задача похожа на See & Spray по closed-loop структуре» или «похожа на робот-такси — домен недостаточен».")
    add_notes(slide, "s36")
    return slide


# ====================================================================
# Q&A (s37)
# ====================================================================

def s37(p):
    """qa — simple farewell. NO photo-hero, NO career cards, minimal."""
    slide = blank(p); set_slide_bg(slide, SURFACE)
    # main course phrase in motif box, centered
    rounded_box(slide, 1.7, 1.7, 9.93, 1.9, fill=WHITE, stroke=GOLD, stroke_w=2.0)
    multiline_box(slide, 2.1, 1.95, 9.13, 1.4, [
        ([("Знать ИИ — значит знать его ", {"size": 30, "bold": True, "color": DEEP}),
          ("границы", {"size": 30, "bold": True, "color": GOLD}),
          (".", {"size": 30, "bold": True, "color": DEEP})], {}),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    # Вопросы?
    text_box(slide, 0.5, 4.0, 12.33, 0.9, "Вопросы?", size=44, bold=True, color=MID,
             align=PP_ALIGN.CENTER)
    text_box(slide, 0.5, 4.95, 12.33, 0.45, "До новых встреч", size=18, italic=True,
             color=DARK_GREY, align=PP_ALIGN.CENTER)
    # three anchors
    rounded_box(slide, 1.7, 5.7, 9.93, 0.85, fill=WHITE, stroke=LIGHT, stroke_w=1.25)
    text_box(slide, 1.95, 5.82, 9.43, 0.6,
             "Две оси карты: применимость × автономия   ·   Семь критериев — один ✗ → СТОП   ·   Три мега-паттерна провалов",
             size=12.5, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    add_notes(slide, "s37")
    return slide


# ====================================================================
# MAIN
# ====================================================================

def main():
    p = setup_pres()
    builders = [
        s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
        s11, s12, s13, s14, s15, s16, s17, s18, s19, s20,
        s21, s22, s23, s24, s25, s26, s27, s28, s29, s30,
        s31, s32, s33, s34, s35, s36, s37,
    ]
    for b in builders:
        b(p)
    assert len(p.slides) == 37, f"expected 37 slides, got {len(p.slides)}"
    p.save(str(OUT))
    print(f"Saved {OUT} with {len(p.slides)} slides")


if __name__ == "__main__":
    main()
