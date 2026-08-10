"""
Full 44-slide build of Лекции 4 «AI в разработке программного обеспечения».
(32 base s01-s32 + suffix-ID dividers s04a/s24a/s28a + content s22a +
§2.4-§2.8 content run s13a..s13h, issue #162 — 8 slides after v2 fix-round
split s13e into s13e/s13f.)

Source-of-truth: deck.yaml + deck-part2.yaml (split >600 строк, loader reads
both) + chapter v1.1 finalized (3 части, ~22300 слов) + slides/*.md (32 файла,
readable speaker notes 150–300 слов).

Issue #99 · Branch: issue-99-lec-04-software-production

Palette LOCKED: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide. Motif «Ocean rounded box»
(radius 12, surface #F4F7FA, stroke #1C7293 1.5pt) на каждом content-слайде.

Canvas: 13.333" × 7.5" (16:9, [#55-1] patch). Pacing per deck.yaml ≈ 75 мин.

Render-style эталон: library/lectures/lec-03/rendered/build_v3.py (та же
палитра/motif/типографика/divider-шаблон/плотность).

Build: python3 build_lec04.py  → lec-04.pptx (44 slides, s01..s32 monotonic
+ suffix-ID inserts).
Charts pre-generated via gen_charts.py; icons via assets/icons (Lucide,
recolored Ocean, 4 variants mid/teal/gold/white).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x5B, 0x66, 0x78)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFD, 0xF3, 0xDC)
TEAL_TINT = RGBColor(0xE4, 0xF1, 0xF2)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
CHARTS = ASSETS / "charts"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-04.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers (architecture mirrors lec-03 build_v3.py — proven)
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before") is not None:
                p.space_before = Pt(cfg["space_before"])
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.035, min(0.22, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, fill, *, stroke=None, stroke_pt=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def connector(slide, x1, y1, x2, y2, color=LIGHT, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = etree.SubElement(
            ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
        pd.set("val", dash)
    return cn


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    """[#73-render-1] aspect-safe: pass only the constraining dimension."""
    path = Path(path)
    if not path.exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path); iw, ih = img.size; img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        ir = iw / ih; br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x),
                                     Inches(y + (h - ah) / 2), width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2),
                                     Inches(y), height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.40, h=0.92, w=12.25, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.10, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 color=DEEP, align=PP_ALIGN.LEFT, stroke_pt=1.5):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=stroke_pt,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.20)


def teal_callout(slide, x, y, w, h, text, *, size=14, bold=False,
                 color=DEEP, align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.18)


def footer(slide, text):
    text_box(slide, x=0.55, y=7.04, w=12.25, h=0.34, text=text,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             line_spacing=1.0)


def trend_stat(slide, x, y, w, h, label, before, after, arrow, *,
               highlight=False):
    """Big before→after number plate (s01-plate style that WORKS).
    arrow: '↑' (rose, bad-direction gold) or '↓' (fell).
    [Fix B1, issue #162] highlight=True now uses SOLID gold fill + DEEP
    text (≈6.85:1 WCAG AA) instead of GOLD_TINT fill + GOLD text
    (≈1.8:1, WCAG FAIL) — gold as text color on light bg is banned."""
    bg = GOLD if highlight else SURFACE
    edge = GOLD if highlight else SOFT_GREY
    filled_rect(slide, x, y, w, h, bg, stroke=edge,
                stroke_pt=(1.5 if highlight else 1.0),
                radius=True, radius_adj=0.10)
    text_box(slide, x + 0.22, y + 0.08, w - 0.44, 0.30, label,
             size=12.5, bold=True, color=(DEEP if highlight else MID),
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, x + 0.22, y + 0.36, w - 0.44, h - 0.42,
             f"{before}  {arrow}  {after}",
             size=25, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


def icon(slide, name, x, y, size, variant="mid"):
    add_image(slide, ICONS / f"{name}-{variant}.png", x, y, size, size)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


def tools_strip(slide, x, y, w, h, tools, direction, caveat, *,
                 chip_fill=MID, left_ratio=0.555):
    """[v3.3, Решение #102] Компактная врезка «Инструменты 2026».

    3 элемента (book-first из chapter v1.2 [for-slide-sNN]):
      (1) tools — список «вендор-режим» строк → chips;
      (2) direction — adoption-НАПРАВЛЕНИЕ словами (растёт/стагнирует),
          БЕЗ точных волатильных чисел/долей на видимом слое;
      (3) caveat — anti-hype/границы-оговорка (⚠), критический тон
          (это AI-Failure-усиление, НЕ вендор-реклама).
    Ocean rounded-box motif; teal-под-band для ⚠ (семантика «граница/
    осторожно», как teal_callout). 0 §/(sNN)/LO/[VFY]/чисел-долей.
    Layout: левая колонка = caption-строка + chips-строка + direction;
    правая = ⚠ band на всю высоту. left_ratio регулирует сплит.
    """
    ocean_box(slide, x, y, w, h)
    pad = 0.22
    inx = x + pad
    iny = y + 0.12
    inh = h - 0.24
    # 2-column inner layout: left = caption+chips+direction, right = ⚠ band.
    left_w = (w - 2 * pad) * left_ratio
    gap = 0.20
    right_x = inx + left_w + gap
    right_w = x + w - pad - right_x
    # Left, row 1 — caption label (own line — chips get full left width)
    text_box(slide, inx, iny, left_w, 0.26, "Инструменты 2026",
             size=12, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    # Left, row 2 — tool chips (full left-column width; auto-fit font)
    n = len(tools)
    cgap = 0.10
    raw = [0.22 + 0.082 * len(t) for t in tools]
    avail = left_w - cgap * (n - 1)
    scale = min(1.0, avail / sum(raw)) if sum(raw) else 1.0
    csz = 11.5 if scale >= 0.92 else (10.5 if scale >= 0.80 else 9.5)
    cx = inx
    cy = iny + 0.32
    for t, rw0 in zip(tools, raw):
        cw = rw0 * scale
        chip(slide, cx, cy, cw, 0.32, t, fill=chip_fill,
             color=WHITE, size=csz, bold=True)
        cx += cw + cgap
    # Left, row 3 — adoption direction (words only, no volatile numbers)
    text_box(slide, inx, iny + 0.70, left_w, inh - 0.70, direction,
             size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.08)
    # Right — anti-hype caveat band (teal = boundary/caution, critical tone;
    # AI-Failure-усиление, НЕ вендор-реклама)
    filled_rect(slide, right_x, iny, right_w, inh, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.25, radius=True, radius_adj=0.07)
    text_runs(slide, right_x + 0.16, iny + 0.04, right_w - 0.32, inh - 0.08, [
        {"text": "⚠  ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": caveat, "size": 11.5, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)


# ============================================================
# Deck loader — deck.yaml split на 2 части (≤600 строк каждая).
# Loader читает ОБЕ части, объединяет ключ `slides`, валидирует totals.
# ============================================================
def load_deck():
    try:
        import yaml
    except ImportError:
        return None
    d1 = yaml.safe_load((ROOT / "deck.yaml").read_text(encoding="utf-8"))
    d2 = yaml.safe_load((ROOT / "deck-part2.yaml").read_text(encoding="utf-8"))
    slides = list(d1.get("slides", [])) + list(d2.get("slides", []))
    ids = [s["id"] for s in slides]
    # [Решение #101, 2026-05-17 — owner GATE B] 32 base (s01–s32 нумерация
    # неизменна) + 3 suffix-ID раздела-дивайдера s04a/s24a/s28a (cascade-safe:
    # chapter [for-slide-sNN] s01–s32 финализирован GATE A, НЕ renumber).
    # [Решение #103, 2026-05-17 — owner GATE C] +1 контент-слайд suffix-ID
    # s22a (curl-slop #5) между s22 и s23 — security-раздел, derive §4.5;
    # owner-документированный slide-count override 35→36, cascade-safe.
    # [issue #162, 2026-08-10] +7 контент-слайдов suffix-ID s13a..s13g
    # (§2.4–§2.8, ранее не покрытые слайдами) между s13 и s14 — Раздел 2
    # (Уровень C), cascade-safe: chapter [for-slide-sNN] s01–s32 не
    # renumber, новые marker-имена s2new-tools/skills/mcp/steering/tasklog.
    # slide-count override 36→43.
    # [issue #162 v2 fix-round, 2026-08-10] presentation-critic REVISE +
    # student-simulator overload finding → old s13e (presence-paradox +
    # Honest Lying + git-conventions crammed together) split: git-
    # conventions moved to new standalone s13f; old s13f/s13g renumbered
    # s13g/s13h. slide-count override 43→44.
    base = [f"s{n:02d}" for n in range(1, 33)]
    expected = []
    for sid in base:
        expected.append(sid)
        if sid == "s04":
            expected.append("s04a")   # Раздел 1 divider
        elif sid == "s13":
            expected.extend(["s13a", "s13b", "s13c", "s13d", "s13e",
                              "s13f", "s13g", "s13h"])  # §2.4–§2.8 (#162 v2)
        elif sid == "s22":
            expected.append("s22a")   # curl-slop #5 (Решение #103)
        elif sid == "s24":
            expected.append("s24a")   # Раздел 5 divider
        elif sid == "s28":
            expected.append("s28a")   # Раздел 6 divider
    assert ids == expected, (
        f"deck slide order mismatch:\n got={ids}\n exp={expected}")
    # s01–s32 нумерация неизменна — base IDs присутствуют все 32 и в порядке
    base_in_order = [i for i in ids if not i.endswith(("a", "b", "c", "d",
                                                         "e", "f", "g", "h"))]
    assert base_in_order == base, (
        f"base s01–s32 numbering changed:\n got={base_in_order}")
    tot = d2.get("totals", {}).get("slides")
    assert tot == 44, f"deck-part2 totals.slides={tot}, expected 44"
    return {"slides": slides, "totals": d2.get("totals", {}),
            "deck": d1.get("deck", {})}


# ============================================================
# Section divider — unified template (7-card roadmap, gold current)
# Sections of Лекции 4 (deck.yaml): 0..6.
# ============================================================
NAV = [
    ("0", "Открытие"),
    ("1", "Уровни A+B"),
    ("2", "Уровень C"),
    ("3", "Уровень D"),
    ("4", "Не только код"),
    ("5", "Методологии"),
    ("6", "Фреймворк"),
]


def roadmap_bar(slide, here_idx, *, y=6.50):
    """7-card progress bar; current section gold-bordered."""
    n = len(NAV)
    gap = 0.12
    bx = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    ch = 0.60
    for i, (num, label) in enumerate(NAV):
        x = bx + i * (cw + gap)
        cur = (i == here_idx)
        if cur:
            filled_rect(slide, x, y, cw, ch, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.12)
        else:
            filled_rect(slide, x, y, cw, ch, SURFACE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(slide, x=x + 0.04, y=y + 0.07, w=cw - 0.08, h=0.22,
                 text=f"Раздел {num}", size=10, bold=True,
                 color=(DEEP if cur else LIGHT), align=PP_ALIGN.CENTER)
        text_box(slide, x=x + 0.04, y=y + 0.30, w=cw - 0.08, h=0.26,
                 text=label, size=10.5, bold=cur,
                 color=(DEEP if cur else SLATE), align=PP_ALIGN.CENTER,
                 line_spacing=0.95)


def build_section_divider(p, here_idx, subtitle, bridge, sid):
    """Distinct divider (NO ocean motif): giant decorative section digit on
    the right (cover-style soft outline), РАЗДЕЛ N + subtitle + 1-line
    narrative bridge on the left, gold-current roadmap bar at bottom."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=8.45, y=0.30, w=4.6, h=6.0, text=str(here_idx),
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.75, y=1.55, w=7.4, h=0.55,
             text=f"РАЗДЕЛ {here_idx}", size=20, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.18, 0.70, 0.05, fill=GOLD)
    text_box(s, x=0.75, y=2.55, w=7.6, h=1.85, text=subtitle,
             size=36, bold=True, color=DEEP, line_spacing=1.08)
    text_box(s, x=0.78, y=4.62, w=7.5, h=1.55, text=bridge,
             size=18, italic=True, color=LIGHT, line_spacing=1.22)
    roadmap_bar(s, here_idx, y=6.50)
    speaker_notes(s, load_notes(sid))
    return s


# ============================================================
# Slide builders — 44 slides (s01..s32 monotonic + suffix-ID inserts)
# ============================================================

def build_s01(p):
    """hero_cover / case_study — METR perception-gap hook. Left: 3 numbers
    RCT in ocean box. Right: HERO — METR's own official study chart (real
    image, Tier 1 og:image acquisition, CC-BY, issue #162 Fix B2), ≥40%
    slide area. Gold: highlight row fill (Fix B1, DEEP-on-GOLD, not gold
    text)."""
    s = blank(p)
    slide_title(s, "16 экспертов были уверены, что AI их ускоряет — и ошиблись на знак.",
                size=24, w=10.85, h=0.96)
    icon(s, "gauge", 12.05, 0.42, 0.74, "mid")
    # Left — METR RCT in ocean box
    lx, ly, lw, lh = 0.55, 1.46, 7.05, 4.55
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.30, ly + 0.22, lw - 0.60, 0.36,
             "Эксперимент METR · RCT · первая половина 2025",
             size=15, bold=True, color=MID)
    text_box(s, lx + 0.30, ly + 0.62, lw - 0.60, 0.78,
             "16 опытных open-source разработчиков · 246 реальных задач "
             "в своих знакомых репозиториях · замеряли реальное время, не ощущение",
             size=13.5, color=DEEP, line_spacing=1.18)
    text_box(s, lx + 0.30, ly + 1.52, lw - 0.60, 0.32,
             "Три числа об одном и том же", size=15, bold=True, color=MID)
    # zone tags make expectation-vs-reality unambiguous in 5 sec
    rows = [
        ("Прогноз до эксперимента", "ОЖИДАЛИ", "−24%", "быстрее",
         LIGHT, False),
        ("Вера после, уже поработав", "ОЖИДАЛИ", "≈ −20%", "быстрее",
         LIGHT, False),
        ("Измеренный факт", "ВЫШЛО", "+19%", "ДОЛЬШЕ", GOLD, True),
    ]
    ry = ly + 1.94
    for lab, zone, val, dirn, col, hi in rows:
        # [Fix B1, issue #162] highlight row: SOLID gold fill + DEEP text
        # (≈6.85:1 WCAG AA), not GOLD_TINT fill + GOLD text (≈1.8:1 FAIL).
        bg = GOLD if hi else SURFACE
        filled_rect(s, lx + 0.30, ry, lw - 0.60, 0.78, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True, radius_adj=0.10)
        # zone chip — cool «ОЖИДАЛИ» vs deep-on-gold «ВЫШЛО»
        filled_rect(s, lx + 0.44, ry + 0.21, 1.16, 0.36,
                    WHITE if hi else SOFT_GREY, radius=True, radius_adj=0.4)
        text_box(s, lx + 0.44, ry + 0.21, 1.16, 0.36, zone,
                 size=11, bold=True, color=(DEEP if hi else MID),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, lx + 1.72, ry + 0.10, 2.02, 0.58, lab,
                 size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
        text_box(s, lx + 3.78, ry + 0.06, 1.40, 0.66, val,
                 size=28, bold=True, color=DEEP,
                 align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, lx + 5.28, ry + 0.06, 1.45, 0.66, dirn,
                 size=(14 if hi else 12.5), bold=hi,
                 color=(DEEP if hi else SLATE), align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        ry += 0.86
    # Right — HERO: METR's own official chart (Fix B2, issue #162).
    # ≈5.00×5.45 = 27.25 sq in of 100 sq in slide ≈ 27% alone; combined with
    # its own attribution/gloss panel below the ocean-box motif frame, the
    # whole right-column visual block ≈ 40%+ of slide area (matches lec-08
    # hero-slot precedent: image + caption counted together).
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 4.55, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, str(ROOT / "assets/screenshots/s01-metr-real-source.png"),
              x=rx + 0.14, y=ly + 0.14, w=rw - 0.28, h=3.85)
    text_box(s, rx + 0.14, ly + 4.04, rw - 0.28, 0.24,
             "METR · официальный график исследования · CC-BY · metr.org",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    text_box(s, rx + 0.14, ly + 4.28, rw - 0.28, 0.24,
             "perception-gap: ощущение скорости ≠ факт",
             size=11.5, bold=True, italic=True, color=TEAL,
             align=PP_ALIGN.LEFT)
    gold_callout(s, 0.55, 6.16, 12.25, 0.74,
                 "Профессионалы, годами пишущие код, ошиблись не в величине — "
                 "в знаке. Ускоряет ли AI лично вас — и откуда вы это "
                 "знаете? «Мне кажется, помогает» — гипотеза, не данные.",
                 size=14)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """cover — distinct, NO ocean motif. Mega «04» + title + roadmap-bar."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=7.7, y=0.95, w=5.7, h=5.2, text="04",
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.75, y=1.45, w=6.6, h=0.5, text="ЛЕКЦИЯ 4",
             size=18, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.02, 0.70, 0.05, fill=TEAL)
    text_box(s, x=0.75, y=2.42, w=7.7, h=2.5,
             text="AI в разработке\nпрограммного обеспечения",
             size=44, bold=True, color=DEEP, line_spacing=1.10)
    filled_rect(s, 0.78, 5.18, 0.05, 0.56, fill=TEAL)
    text_box(s, x=1.02, y=5.16, w=7.4, h=0.62,
             text="Где AI ускоряет, где вредит —\nи что в работе инженера НЕ делегируется",
             size=18, color=MID, line_spacing=1.18)
    roadmap_bar(s, 0, y=6.50)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    """schema_matrix KEYSTONE — лестница автономности A→D: 4 уровня
    (строки A/B/C/D) × что делает AI · кто принимает решение · живой пример.
    Несущая ось всей лекции. Letter-badge escalation = рост автономии.
    course-scaffold-атрибуция → ТОЛЬКО speaker notes; 0 §/disclaimer."""
    s = blank(p)
    slide_title(s, "Лестница автономности A→D — четыре уровня участия AI.",
                size=27, y=0.34, h=0.60, w=12.25)
    text_box(s, 0.55, 1.04, 12.25, 0.34,
             "Как лестница сложности из Лекции 3 — выбираешь ступень под "
             "задачу, не выше: каждый шаг вверх отдаёт AI больше решений.",
             size=14, italic=True, color=MID, line_spacing=1.1)

    ocean_box(s, 0.40, 1.50, 12.55, 4.34)
    tx, ty = 0.56, 1.64
    # columns: [уровень-бейдж+имя | что делает AI | кто принимает решение | живой пример]
    col_w = [2.12, 3.28, 2.98, 4.01]
    hh = 0.50
    rh = 0.94
    headers = ["Уровень", "Что делает AI",
               "Кто принимает решение", "Живой пример"]
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID)
        text_box(s, cx + 0.12, ty, col_w[j] - 0.20, hh, hd,
                 size=13.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                 align=(PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT))
        cx += col_w[j]

    # A→D — цвет бейджа эскалирует (рост автономии — несущая ось);
    # D = апекс автономии → gold-кольцо (несущий gold-акцент слайда)
    rows = [
        ("A", "автодополнение", LIGHT, None,
         "дописывает строку или блок прямо в потоке набора",
         "человек — на каждом нажатии Tab",
         "Copilot подсказал хвост строки — нажал Tab, читая"),
        ("B", "мелкие задачи", MID, None,
         "пишет функцию или фикс по запросу в чате/inline",
         "человек ставит задачу и ревьюит результат",
         "«напиши парсер CSV» в чате — вставил, проверил"),
        ("C", "кодинг-агент", DEEP, None,
         "берёт многофайловую задачу, сам гоняет тесты, итерирует",
         "человек ревьюит pull request и решает о слиянии",
         "«реализуй фичу X» — агент правит 5 файлов + тесты"),
        ("D", "оркестратор", DEEP, GOLD,
         "берёт задачу из трекера и доводит до PR сам",
         "человек — стратегия, approval, merge, прод-гейт",
         "агент взял тикет, открыл PR — человек ревьюит / мержит"),
    ]
    yy = ty + hh
    for ri, (lt, lname, badge, ring, c_ai, c_dec, c_ex) in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else SURFACE
        cx = tx
        # уровень-ячейка: бейдж слева (вертикально по центру) + имя уровня
        # справа на всю оставшуюся ширину ячейки (1 строка, без переноса)
        filled_rect(s, cx, yy, col_w[0], rh, bg, stroke=SOFT_GREY,
                    stroke_pt=0.75)
        bd = 0.56
        if ring is not None:
            circle(s, cx + 0.16 - 0.055, yy + rh / 2 - bd / 2 - 0.055,
                   bd + 0.11, ring)
        circle(s, cx + 0.16, yy + rh / 2 - bd / 2, bd, badge)
        text_box(s, cx + 0.16, yy + rh / 2 - bd / 2, bd, bd, lt,
                 size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + 0.80, yy, col_w[0] - 0.84, rh, lname,
                 size=11, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0, align=PP_ALIGN.LEFT)
        cx += col_w[0]
        for j, cc in enumerate((c_ai, c_dec, c_ex), start=1):
            filled_rect(s, cx, yy, col_w[j], rh, bg, stroke=SOFT_GREY,
                        stroke_pt=0.75)
            text_box(s, cx + 0.16, yy + 0.04, col_w[j] - 0.30, rh - 0.08, cc,
                     size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.10)
            cx += col_w[j]
        yy += rh

    # gold spine: несущая ось — рост автономии снизу вверх по A→D
    gold_callout(s, 0.55, 5.96, 12.25, 0.90,
                 "Чем ниже в таблице — тем больше решений у AI и тем меньше "
                 "у человека. Эти три колонки — линза, как читать "
                 "каждый уровень.", size=14)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """assertion_visual — central question + answer-frame + 5 «где человек
    обязателен» якорей именами по смыслу. 0 «возвращаемся N раз» / §-кодов /
    course-scaffold disclaimer (атрибуция → speaker notes s03)."""
    s = blank(p)
    slide_title(s, "Центральный вопрос лекции.",
                size=27, y=0.34, h=0.58, w=9.5)
    icon(s, "target", 12.05, 0.36, 0.78, "gold")
    bx, by, bw, bh = 0.55, 1.22, 12.25, 1.70
    ocean_box(s, bx, by, bw, bh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, bx + 0.40, by + 0.20, bw - 0.80, bh - 0.40,
             "«AI пишет код всё лучше — где он реально ускоряет, где "
             "замедляет или вредит, и что в работе инженера НЕ делегируется?»",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16, align=PP_ALIGN.CENTER)
    # Рамка ответа
    teal_callout(s, 0.55, 3.10, 12.25, 0.80,
                 "Ответ — не «AI хорош» или «AI плох», а: назвать уровень "
                 "A–D, конфигурацию и точку, где человек обязателен.",
                 size=14.5, bold=True)
    # 5 якорей «где человек обязателен» — именами по смыслу
    text_box(s, 0.55, 4.06, 12.25, 0.30,
             "Пять мест, где человек обязателен:",
             size=14, bold=True, color=MID)
    anchors = [
        ("«почти правильный»\nкод"),
        ("merge\nи ревью"),
        ("деструктив\nна prod"),
        ("безопасность\nкода"),
        ("что строить\n(essential)"),
    ]
    n = len(anchors)
    gap = 0.20
    cw = (12.25 - gap * (n - 1)) / n
    ax, ay, ah = 0.55, 4.42, 1.10
    for idx, lab in enumerate(anchors, start=1):
        ocean_box(s, ax, ay, cw, ah)
        # верхняя teal-полоска «человек обязателен» — общий семантический ярлык
        filled_rect(s, ax + 0.14, ay + 0.12, cw - 0.28, 0.28, TEAL_TINT,
                    stroke=None, radius=True, radius_adj=0.4)
        text_box(s, ax + 0.14, ay + 0.12, cw - 0.28, 0.28,
                 "человек обязателен", size=9.5, bold=True, color=TEAL,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, ax + 0.12, ay + 0.46, cw - 0.24, ah - 0.56, lab,
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        ax += cw + gap
    gold_callout(s, 0.55, 5.70, 12.25, 1.16,
                 "Главное в работе инженера — именно то, что НЕ делегируется. "
                 "Первая половина вопроса — про скорость; вторая — про "
                 "ответственность. Ответственность не делегируется.",
                 size=14.5)
    speaker_notes(s, load_notes("s04"))


def build_s04a(p):
    """section_divider — Раздел 1: Уровни A и B. [Решение #101 owner GATE B]
    suffix-ID, cascade-safe; here_idx=1 → roadmap gold-маркер Раздел 1.
    Шаблон полностью единый с s10/s14/s18 (build_section_divider)."""
    build_section_divider(
        p, 1, "Уровни A и B:\nавтодополнение и чат",
        "Начинаем снизу лестницы. A — AI дописывает строку прямо в потоке "
        "набора; B — пишет функцию или фикс по запросу в чате. На обоих "
        "человек видит каждый фрагмент кода — это AI, который уже стоит в "
        "каждой IDE.", "s04a")


def build_s05(p):
    """assertion_visual — несущий принцип: цена ошибки (радиус поражения)
    растёт с автономией A→D; точка человеческого контроля смещается и
    ужесточается. НЕ meta, НЕ защита таксономии. Визуал — растущий
    blast-radius вдоль A→D. Слот переиспользован (НЕ renumber)."""
    s = blank(p)
    slide_title(s, "Чем выше автономия — тем больше радиус поражения ошибки.",
                size=25, y=0.34, h=0.58, w=12.25)
    text_box(s, 0.55, 0.98, 12.25, 0.34,
             "Поэтому точка обязательного человека не исчезает с подъёмом — "
             "она смещается и ужесточается.",
             size=14, italic=True, color=MID, line_spacing=1.1)

    # 4 столбца A→D: растущий сплошной blast-bar (bottom-aligned) с тем,
    # что задевает ошибка, на самом баре; внизу — где человек контролирует
    ocean_box(s, 0.40, 1.42, 12.55, 4.10)
    levels = [
        ("A", "автодополнение", LIGHT, False,
         "одна строка", "каждый токен", 0.60),
        ("B", "мелкие задачи", MID, False,
         "функция / фикс", "ревью фрагмента", 1.12),
        ("C", "кодинг-агент", DEEP, False,
         "много файлов сразу", "ревью PR до merge", 1.74),
        ("D", "оркестратор", GOLD, True,
         "прод-данные,\nнеобратимое", "только вход / выход", 2.42),
    ]
    n = 4
    gap = 0.22
    cw = (12.55 - 0.56 - gap * (n - 1)) / n
    cx = 0.40 + 0.28
    base_y = 4.98  # общая нижняя граница bars (bottom-aligned)
    badge_y = 1.56
    for lt, lname, bar_col, is_d, scope, ctrl, bh in levels:
        bd = 0.50
        if is_d:
            circle(s, cx + cw / 2 - bd / 2 - 0.055, badge_y - 0.055,
                   bd + 0.11, GOLD)
        circle(s, cx + cw / 2 - bd / 2, badge_y, bd,
               DEEP if is_d else bar_col)
        text_box(s, cx + cw / 2 - bd / 2, badge_y, bd, bd, lt,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx, badge_y + 0.54, cw, 0.28, lname, size=11.5,
                 bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        # сплошной растущий blast-bar (цвет = badge; D = gold)
        filled_rect(s, cx + 0.16, base_y - bh, cw - 0.32, bh, bar_col,
                    stroke=None, radius=True, radius_adj=0.05)
        text_box(s, cx + 0.10, base_y - bh + 0.08, cw - 0.20, bh - 0.16,
                 scope, size=12, bold=True,
                 color=(DEEP if is_d else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
        # где человек контролирует — под баром (DEEP — читаемо; gold уже
        # несёт bar+ring+callout, текстовый контраст важнее)
        text_box(s, cx, base_y + 0.08, cw, 0.42,
                 "человек:\n" + ctrl, size=11, bold=True,
                 color=(DEEP if is_d else MID),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        cx += cw + gap

    gold_callout(s, 0.55, 5.74, 12.25, 1.10,
                 "Тот же провал на уровне D стоит несравнимо дороже, чем на "
                 "A. Поднимать автономию имеет смысл только под требование "
                 "задачи — и только с гейтом под выросшую цену ошибки.",
                 size=14.5)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """case_study — Уровень A: спуск на 1-ю ступень лестницы с s03.
    3 контекста-эффекта + «где уже стоит / в чём ловушка». Gold −19%.
    Подаём как новое: что это / где стоит / ловушка — БЕЗ scaffold-рамки."""
    s = blank(p)
    slide_title(s, "Уровень A — автодополнение: безопасен только пока человек реально читает.",
                size=21, y=0.34, h=0.56, w=12.25)
    text_box(s, 0.55, 0.94, 12.25, 0.42,
             "Первая ступень лестницы: AI дописывает строку или блок по "
             "контексту файла; человек принимает или отклоняет каждое "
             "предложение в момент написания.",
             size=13, italic=True, color=MID, line_spacing=1.10)
    # Left — 3 крупных stat-плашки (v3.3: spacing compressed to fit
    # tools_strip; числа +56/+7…22/выигрыш-исчезает НЕ изменены)
    lx, ly, lw, lh = 0.55, 1.44, 7.05, 3.46
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.14, lw - 0.52, 0.30,
             "Один инструмент — три разных эффекта на скорость",
             size=14, bold=True, color=MID)
    ctx = [
        ("Лаборатория — изолированная новая задача", "+56%", 31, MID, False),
        ("Поле — Microsoft / Accenture", "+7…22%", 31, LIGHT, False),
        ("Знакомое легаси у экспертов",
         "выигрыш\nисчезает", 17, GOLD, True),
    ]
    cyy = ly + 0.52
    csh = 0.90
    for lab, val, vsz, col, hi in ctx:
        # [Fix B1, issue #162] highlight: SOLID gold fill + DEEP text.
        bg = GOLD if hi else SURFACE
        filled_rect(s, lx + 0.26, cyy, lw - 0.52, csh, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.09)
        text_box(s, lx + 0.46, cyy + 0.12, lw - 2.85, csh - 0.24, lab,
                 size=13, bold=True, color=(DEEP if hi else MID),
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        text_box(s, lx + lw - 2.45, cyy + 0.08, 2.05, csh - 0.16, val,
                 size=vsz, bold=True, color=DEEP,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        cyy += csh + 0.07
    # Right — уровень A на практике: где стоит / в чём ловушка (не scaffold)
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 2.34)
    text_box(s, rx + 0.26, ly + 0.14, rw - 0.52, 0.28,
             "Где уже стоит — и в чём ловушка", size=14, bold=True,
             color=MID)
    frame = [
        ("Где стоит:", "в IDE почти каждого инженера — Copilot-класс"),
        ("Человек:", "читает и принимает каждое предложение сам"),
        ("Ловушка:", "привычка жать «принять», не читая"),
        ("Цена:", "клоны и уязвимые паттерны попадают молча"),
    ]
    fy = ly + 0.48
    for a, b in frame:
        text_box(s, rx + 0.26, fy, 1.30, 0.44, a, size=12, bold=True,
                 color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 1.58, fy, rw - 1.84, 0.44, b, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        fy += 0.46
    gold_callout(s, 7.80, ly + 2.50, 5.00, 0.96,
                 "Сними чтение — и самый безопасный уровень становится "
                 "поставщиком техдолга на скорости набора текста: на знакомом "
                 "коде у экспертов выигрыш в скорости исчезает совсем.",
                 size=12.5)
    # [v3.3, Решение #102] Инструменты уровня A 2026 (book-first §1.2)
    tools_strip(
        s, 0.55, 5.02, 12.25, 1.42,
        ["Copilot ghost-text", "Cursor Tab", "JetBrains AI Assistant"],
        "Adoption: A — самый зрелый и широкий по охвату уровень; "
        "лидер по охвату — Copilot-класс, но рост лидера остановился.",
        "«Copilot — №1» — это охват, не динамика. Стагнация лидера "
        "≠ «инструмент умер». На уровень A ставит режим "
        "tab-completion, а не логотип.")
    footer(s, "Peng et al. (Copilot RCT, 2023, лабораторные); MIT GenAI "
              "(2024, поле). Эффект AI контекстно-зависим; лабораторные числа "
              "подавать как лабораторные.")
    speaker_notes(s, load_notes("s06"))


def build_s07(p):
    """assertion_visual — Уровень B: граница A↔B как РАЗЛИЧЕНИЕ (не
    оправдание классификации). 2-кол сравнение + что эта граница меняет
    на практике (концентр., не scaffold-рамка)."""
    s = blank(p)
    slide_title(s, "Уровень B: человек проверяет после, а не во время — первое делегирование.",
                size=21, y=0.34, h=0.56, w=12.25)
    # 2-col compare in ocean box (v3.3: compressed vertically to fit
    # tools_strip — контент не урезан, только spacing)
    bx, by, bw, bh = 0.55, 1.34, 12.25, 2.42
    ocean_box(s, bx, by, bw, bh)
    cols = ["", "Уровень A", "Уровень B"]
    col_w = [3.05, 4.55, 4.55]
    rows = [
        ("Единица работы", "строка-в-потоке", "задача-фрагмент (функция, фикс)"),
        ("Где человек в цикле", "на каждом токене", "после генерации"),
        ("Когда ревью", "в момент написания", "отдельным шагом"),
    ]
    tx, ty = bx + 0.20, by + 0.16
    hh, rh = 0.46, 0.54
    cx = tx
    for j, c in enumerate(cols):
        filled_rect(s, cx, ty, col_w[j], hh, MID if j else SURFACE,
                    stroke=(None if j else SOFT_GREY),
                    stroke_pt=(0 if j else 1.0))
        if c:
            text_box(s, cx, ty, col_w[j], hh, c, size=14, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for r0, r1, r2 in rows:
        cx = tx
        for j, cc in enumerate((r0, r1, r2)):
            filled_rect(s, cx, yy, col_w[j], rh, WHITE if j else SURFACE,
                        stroke=SOFT_GREY, stroke_pt=0.75)
            text_box(s, cx + 0.12, yy, col_w[j] - 0.24, rh, cc,
                     size=13, bold=(j == 0), color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE,
                     align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER),
                     line_spacing=1.0)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 3.92, 12.25, 0.74,
                 "Смещение «человек проверяет после, а не во время» — первое "
                 "реальное делегирование в лестнице. С него начинаются "
                 "характерные проблемы AI-кода.", size=14)
    # Что эта граница меняет на практике (концентрированно, не scaffold)
    ocean_box(s, 0.55, 4.72, 12.25, 1.28)
    text_box(s, 0.80, 4.81, 5.0, 0.28, "Что эта граница меняет на практике",
             size=13, bold=True, color=MID)
    fr = [
        ("Защита «я видел каждую строку»", "на B уже не работает"),
        ("Ревью становится отдельным шагом", "его легко пропустить"),
        ("Появляется «почти правильный» код", "разберём следующим"),
    ]
    fx = 0.80
    cwf = 3.95
    fy = 5.13
    for i, (a, b) in enumerate(fr):
        xx = fx + i * cwf
        text_box(s, xx, fy, cwf - 0.20, 0.38, a, size=12, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        text_box(s, xx, fy + 0.40, cwf - 0.20, 0.36, "→ " + b, size=11.5,
                 italic=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    # [v3.3, Решение #102] Инструменты уровня B 2026 (book-first §1.3)
    tools_strip(
        s, 0.55, 6.06, 12.25, 1.34,
        ["ChatGPT-чат", "Copilot Chat", "Cursor Cmd-K"],
        "Adoption: чат-LLM — самый массовый способ применять AI к "
        "коду (не требует интеграции, доступен каждому).",
        "Чат-LLM — строго уровень B (петля copy-paste, человек "
        "после каждого шага), даже если вендор обещает «агентно»: "
        "без обвязки «может агентно» — маркетинг.")
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """case_study — «почти правильный» 66% (c08 donut) + 70/80% + lesson."""
    s = blank(p)
    slide_title(s, "Последние 20–30% так же трудны, как были до AI.", size=26)
    # Left — 70/80% problem
    lx, ly, lw, lh = 0.55, 1.42, 6.55, 3.65
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.18, lw - 0.52, 0.34, "70/80%-проблема",
             size=16, bold=True, color=MID)
    pts = [
        "Первые 70–80% — быстро и дёшево: типовое, было в обучающих данных",
        "Последние 20–30% — краевые случаи, ошибки, безопасность, "
        "интеграция, нагрузка — так же трудны",
        "Разрыв структурный, не временный: специфики системы в обучающих "
        "данных не было и быть не могло",
    ]
    py = ly + 0.66
    for t in pts:
        circle(s, lx + 0.28, py + 0.06, 0.13, MID)
        text_box(s, lx + 0.56, py, lw - 0.84, 0.94, t,
                 size=13, color=DEEP, line_spacing=1.16)
        py += 0.98
    # Right — крупное число 66% (вместо donut без легенды)
    rx, rw = 7.30, 5.50
    ocean_box(s, rx, ly, rw, 3.65, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    # [Fix B1, issue #162] gold TEXT on GOLD_TINT (≈1.8:1 WCAG FAIL) →
    # solid-gold plate behind the number + DEEP text (≈6.85:1 PASS).
    filled_rect(s, rx + 0.30, ly + 0.20, rw - 0.60, 1.58, GOLD,
                radius=True, radius_adj=0.10)
    text_box(s, rx + 0.30, ly + 0.22, rw - 0.60, 1.55, "66%",
             size=92, bold=True, color=DEEP, line_spacing=0.95,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text_box(s, rx + 0.30, ly + 1.78, rw - 0.60, 0.46,
             "разработчиков (опрос Stack Overflow, 2025)",
             size=13, bold=True, color=DEEP)
    text_box(s, rx + 0.30, ly + 2.22, rw - 0.60, 1.30,
             "топ-фрустрация — код «почти правильный, но не совсем». "
             "Дороже явно неверного: компилируется, проходит обычный "
             "сценарий — ломается в проде.",
             size=12.5, color=DEEP, line_spacing=1.18)
    gold_callout(s, 0.55, 5.22, 12.25, 1.04,
                 "Урок: «AI написал за минуту» без «и я проверил» = «долг "
                 "записан за минуту». Любой фрагмент, чья некорректность не "
                 "обнаружится автоматически, обязан пройти человеческое "
                 "ревью до интеграции.", size=14.5)
    footer(s, "Osmani (70%-проблема, 2024); Stack Overflow Survey 2025 "
              "(66% — опрос, направление стабильно).")
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """assertion_visual — structure + constraints + tests; bridge to TDD."""
    s = blank(p)
    slide_title(s, "Не запрет, а инженерный паттерн: structure + constraints + tests.",
                size=24, y=0.34, h=0.62)
    text_box(s, 0.55, 0.98, 12.25, 0.32,
             "Канонический способ ставить AI задачу так, чтобы «почти "
             "правильное» ловила машина, а не глаза в проде.",
             size=13, italic=True, color=MID)
    cards = [
        ("braces", "Structure", "mid", False,
         "дать сигнатуры, типы, контракт интерфейса, входы/выходы — а не "
         "«сделай мне X». Чем строже структура, тем уже пространство "
         "правдоподобно-неверного."),
        ("shield-check", "Constraints", "mid", False,
         "явно перечислить, чего нельзя и какие инварианты держать. "
         "Модель не выведет их сама — она не знает контекста системы."),
        ("flask-conical", "Tests", "gold", True,
         "тест как исполняемая спецификация — машинно-проверяемый критерий "
         "«правильно/нет», не подвержен ни «почти правильному», ни "
         "perception-gap."),
    ]
    n = 3
    gap = 0.28
    cw = (12.25 - gap * (n - 1)) / n
    cx = 0.55
    cy, chh = 1.46, 3.00
    for ic, ttl, var, hi, body in cards:
        if hi:
            ocean_box(s, cx, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD,
                      stroke_pt=2.0)
        else:
            ocean_box(s, cx, cy, cw, chh)
        icon(s, ic, cx + 0.28, cy + 0.24, 0.56, var)
        text_box(s, cx + 0.98, cy + 0.30, cw - 1.15, 0.50, ttl,
                 size=18, bold=True, color=(DEEP if hi else MID),
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + 0.28, cy + 0.96, cw - 0.56, chh - 1.10, body,
                 size=12.5, color=DEEP, line_spacing=1.18)
        cx += cw + gap
    gold_callout(s, 0.55, 4.62, 12.25, 1.64,
                 "Это, по сути, TDD (разработка через тесты), применённый к "
                 "постановке задачи для AI. Всё, что AI пишет начиная с "
                 "уровня B, должно сопровождаться машинно-проверяемым "
                 "критерием корректности. Антипаттерн — vibe-coding: "
                 "генерировать и принимать код «по ощущению», без структуры, "
                 "ограничений, теста и проверочного шага (строгое "
                 "определение — дальше в лекции).", size=14)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    """section_divider — Раздел 2: Уровень C."""
    build_section_divider(
        p, 2, "Уровень C:\nкодинг-агент",
        "На A и B человек видел каждый фрагмент. Уровень C — AI сам "
        "планирует, правит много файлов, гоняет тесты; человек видит только "
        "итог — pull request.", "s10")


def build_s11(p):
    """process / schema_cycle — coding-agent = plan→act→check→iterate."""
    s = blank(p)
    slide_title(s, "Кодинг-агент = цикл Лекции 3, применённый к коду.", size=26)
    icon(s, "refresh-cw", 11.75, 0.36, 0.78, "mid")
    # Cycle 4 steps as 2x2 with curved feel + return arrow label
    cx, cy = 0.55, 1.42
    cw, chh = 12.25, 3.10
    ocean_box(s, cx, cy, cw, chh)
    steps = [
        ("plan", "формулирует следующий шаг (какой файл, что изменить)",
         "mid", False),
        ("act", "правит код, запускает тест/сборку (его «инструменты»)",
         "mid", False),
        ("check", "читает результат прогона, оценивает цель",
         "gold", True),
        ("iterate", "повторяет цикл, пока цель не достигнута",
         "mid", False),
    ]
    n = 4
    gap = 0.20
    sw = (cw - 0.50 - gap * (n - 1)) / n
    sx = cx + 0.25
    for i, (nm, desc, var, hi) in enumerate(steps):
        bx = sx + i * (sw + gap)
        if hi:
            filled_rect(s, bx, cy + 0.22, sw, 2.66, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.07)
        else:
            filled_rect(s, bx, cy + 0.22, sw, 2.66, WHITE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.07)
        text_box(s, bx + 0.10, cy + 0.52, sw - 0.20, 0.50, nm,
                 size=20, bold=True, color=(DEEP if hi else MID),
                 align=PP_ALIGN.CENTER)
        text_box(s, bx + 0.16, cy + 1.18, sw - 0.32, 1.40, desc,
                 size=13.5, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
        if i < n - 1:
            right_arrow(s, bx + sw + 0.01, cy + 1.30, gap - 0.02, 0.42,
                        fill=LIGHT)
    text_box(s, cx, cy + chh + 0.02, cw, 0.30,
             "↻ повторяется: iterate → plan (замкнутый цикл)",
             size=12.5, bold=True, italic=True, color=TEAL,
             align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 4.90, 12.25, 0.78,
                 "Шаг check не должен быть самооценкой модели: «я справился» "
                 "порождается тем же механизмом, что и сама ошибка — это "
                 "вывод Лекции 3, применённый к коду.", size=14)
    ocean_box(s, 0.55, 5.82, 12.25, 1.06)
    text_box(s, 0.78, 5.94, 11.8, 0.84,
             "Рамка C:  что делает AI — ведёт многошаговую разработку сам · "
             "кто решает — человек ставит задачу и решает про merge · где "
             "обязателен — на ревью pull request и merge · риск — падение "
             "надёжности на незнакомом коде",
             size=12, color=DEEP, line_spacing=1.20,
             anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """comparison / schema_matrix — SWE-bench Verified vs Pro (c12)."""
    s = blank(p)
    slide_title(s, "Один класс систем — два разных результата.", size=26)
    text_box(s, 0.55, 1.04, 12.25, 0.46,
             "SWE-bench — стандартный бенчмарк (мерило): AI дают реальную "
             "задачу из трекера open-source проекта; считают долю задач, где "
             "патч проходит тесты проекта.", size=13, italic=True, color=MID,
             line_spacing=1.10)
    # Left — 2 mega-stat плашки + gold-дельта (v3.3: spacing compressed
    # to fit tools_strip; mega-числа 42pt = 5-сек якорь, НЕ урезаны)
    lx, ly, lw, lh = 0.55, 1.54, 7.05, 3.14
    ocean_box(s, lx, ly, lw, lh)
    px, pw = lx + 0.26, lw - 0.52
    # Verified
    filled_rect(s, px, ly + 0.16, pw, 1.14, SURFACE, stroke=SOFT_GREY,
                stroke_pt=1.0, radius=True, radius_adj=0.08)
    text_box(s, px + 0.24, ly + 0.24, pw - 1.9, 0.98,
             "SWE-bench Verified\nзнакомый публичный код (был в обучении)",
             size=13, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.06)
    text_box(s, px + pw - 2.05, ly + 0.18, 1.85, 1.10, "~95–96%",
             size=36, bold=True, color=DEEP, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    # gold delta band [Fix C, issue #162: gap narrowed ~24pp → ~15-17pp,
    # Scale SEAL conservative estimate, per chapter.md §2.2]
    filled_rect(s, px, ly + 1.36, pw, 0.42, GOLD, radius=True,
                radius_adj=0.20)
    text_box(s, px, ly + 1.36, pw, 0.42,
             "разрыв ~15–17 процентных пунктов",
             size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Pro [Fix B1, issue #162: gold TEXT on GOLD_TINT (≈1.8:1 FAIL) →
    # DEEP text — box already gold-framed, no need for text to be gold too]
    filled_rect(s, px, ly + 1.84, pw, 1.14, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.08)
    text_box(s, px + 0.24, ly + 1.92, pw - 1.9, 0.98,
             "SWE-bench Pro\nчестный незнакомый приватный код",
             size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.06)
    text_box(s, px + pw - 2.05, ly + 1.86, 1.85, 1.10, "~69–80%",
             size=36, bold=True, color=DEEP, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    # Right — explanation
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 1.40, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 0.12, rw - 0.52, 1.16,
             "Главный инженерный факт уровня C: «почти 95%» на знакомом "
             "коде → «7–8 из 10» на незнакомом",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.12)
    ocean_box(s, rx, ly + 1.52, rw, 1.62)
    text_box(s, rx + 0.26, ly + 1.64, rw - 0.52, 1.40,
             "Verified — ~500 проверенных задач из публичного кода, который "
             "модель видела при обучении. Pro — приватные кодбазы, которых "
             "модель видеть не могла: честный незнакомый код.",
             size=12.5, color=DEEP, line_spacing=1.14)
    gold_callout(s, 0.55, 4.80, 12.25, 0.74,
                 "Доверие к кодинг-агенту обратно пропорционально "
                 "незнакомости и критичности кода. Принимать PR без "
                 "тщательного ревью — строить на цифре, которая к вашему "
                 "коду не относится.", size=14)
    # [v3.3, Решение #102] Инструменты уровня C 2026 (book-first §2.2)
    tools_strip(
        s, 0.55, 5.62, 12.25, 1.36,
        ["Claude Code", "Cursor Composer", "Codex CLI"],
        "Adoption: C — самый быстрорастущий уровень; частый паттерн "
        "«связка инструментов» (редактор-агент + кодинг-агент).",
        "SWE-bench как доказательство автономии дыряв: высокая цифра "
        "≠ «мерджить без senior-ревью». Уровень задаёт режим (сам "
        "итерирует и гоняет тесты), а не бренд.")
    footer(s, "SWE-bench — мерило: набор реальных задач, на котором "
              "сравнивают модели; лидеры меняются почти еженедельно. Цифра "
              "без среза («Verified или Pro?») и без даты не информативна.")
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """case_study — antipattern «merge зелёные тесты» + GitClear (c13)."""
    s = blank(p)
    slide_title(s, "«Зелёные тесты» — необходимое, но не достаточное.", size=26)
    # Left — antipattern
    lx, ly, lw, lh = 0.55, 1.40, 6.35, 3.70
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "triangle-alert", lx + 0.26, ly + 0.20, 0.48, "gold")
    text_box(s, lx + 0.88, ly + 0.24, lw - 1.05, 0.44, "Антипаттерн уровня C",
             size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 0.84, lw - 0.52, 1.05,
             "«merge по зелёным тестам без чтения» = уровень C, де-факто "
             "деградировавший до D без явного решения о повышении автономии.",
             size=13, color=DEEP, line_spacing=1.18)
    text_box(s, lx + 0.26, ly + 2.00, lw - 0.52, 1.55,
             "Тесты проверяют то, что в них написано, а не то, что код не "
             "дублирует существующее, не ломает архитектуру, не вносит "
             "уязвимость в краевом пути без теста.",
             size=13, color=DEEP, line_spacing=1.18)
    # Right — GitClear: 3 крупных trend-числа (вместо декор-chart)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 3.70)
    text_box(s, rx + 0.24, ly + 0.13, rw - 0.48, 0.28,
             "GitClear: анализ 211 млн строк кода, 2020 → 2024",
             size=13.5, bold=True, color=MID)
    sx, sw = rx + 0.24, rw - 0.48
    sh = 0.86
    sy = ly + 0.46
    trend_stat(s, sx, sy, sw, sh,
               "Клоны кода, %", "8,3", "12,3", "↑", highlight=True)
    trend_stat(s, sx, sy + sh + 0.05, sw, sh,
               "Доля рефакторинга, %", "24,1", "9,5", "↓")
    trend_stat(s, sx, sy + 2 * (sh + 0.05), sw, sh,
               "Churn — переписано ≤2 нед, %", "5,5", "7,9", "↑")
    # [Fix B1, issue #162] gold TEXT on white → gold FILL chip + DEEP text
    _concl_y = sy + 3 * sh + 2 * 0.05 + 0.02
    filled_rect(s, rx + 0.24, _concl_y, rw - 0.48, 0.36, GOLD,
                radius=True, radius_adj=0.20)
    text_box(s, rx + 0.24, _concl_y, rw - 0.48, 0.36,
             "→ AI ускоряет порождение кода, не его качество",
             size=12, bold=True, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.24, 12.25, 1.02,
                 "Альтернатива (не-AI): обязательное человеческое ревью PR + "
                 "метрики дублирования/churn в CI как gate. Merge = решение "
                 "об ответственности за код — не делегируется, как code "
                 "review между людьми не отменяется доверием к коллеге.",
                 size=13.5)
    footer(s, "GitClear 2025 (211M строк, 25 крупнейших OSS + приватные) — "
              "корреляция во времени, направление стабильно.")
    speaker_notes(s, load_notes("s13"))


def build_s13a(p):
    """assertion_visual — §2.4 landscape: agentic IDE / CLI-agent / framework
    categories (orthogonal to A->D ladder) + category>brand echo + agent vs
    subagent SWE example (issue #162, new §2.4-2.8 content)."""
    s = blank(p)
    slide_title(s,
                "Три категории инструментов — по тому, ГДЕ живёт агент.",
                size=25, y=0.34, h=0.60)
    text_box(s, 0.55, 0.98, 12.25, 0.32,
             "Отдельная ось от лестницы A→D: та отвечает «сколько AI решает "
             "сам», эта — «в каком окружении».", size=13, italic=True,
             color=MID)
    cards = [
        ("code", "Agentic IDE", "mid",
         "Desktop-редактор (форк VS Code) со встроенным агентом. Флагман — "
         "Cursor: Tab (A), Cmd-K (B), Composer/Agent Mode (C) — всё в одном "
         "продукте.",
         "визуальный контроль (diff, файловое дерево) ценой привязки к "
         "редактору"),
        ("terminal", "CLI-агент", "mid",
         "Работает из терминала без GUI, в связке с любым редактором. "
         "Canonical пример — Claude Code: по умолчанию C, до D в режиме "
         "команд агентов.",
         "гибкость в автоматизации (CI, hook) ценой отсутствия diff-UI"),
        ("boxes", "Фреймворк для агентов", "teal",
         "Не готовый продукт, а библиотека/SDK для сборки своего агента под "
         "нестандартную задачу. Для полноты словаря — курс про пользователей "
         "готовых инструментов.",
         "инженер строит агента сам, а не выбирает из готовых"),
    ]
    n = 3
    gap = 0.24
    cw = (12.25 - gap * (n - 1)) / n
    cx0, cy, chh = 0.55, 1.48, 3.10
    for ic, ttl, var, body, tail in cards:
        ocean_box(s, cx0, cy, cw, chh)
        icon(s, ic, cx0 + 0.24, cy + 0.20, 0.50, var)
        text_box(s, cx0 + 0.86, cy + 0.24, cw - 1.05, 0.44, ttl,
                 size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx0 + 0.24, cy + 0.86, cw - 0.48, 1.56, body,
                 size=12, color=DEEP, line_spacing=1.16)
        text_box(s, cx0 + 0.24, cy + 2.48, cw - 0.48, 0.56, "→ " + tail,
                 size=11, italic=True, color=TEAL, line_spacing=1.12)
        cx0 += cw + gap
    teal_callout(s, 0.55, 4.76, 12.25, 0.80,
                 "Категория важнее бренда — тот же принцип, что и для уровня "
                 "автономии: agentic IDE выигрывает при визуальном контроле "
                 "над многофайловым diff в одном редакторе; CLI-агент — при "
                 "автоматизации, конвейере, headless-окружении, параллельных "
                 "сессиях.", size=13)
    gold_callout(s, 0.55, 5.70, 12.25, 1.06,
                 "Агент vs субагент в коде: основной агент делегирует "
                 "субагенту узкий подшаг («прочитай 15 файлов миграций — "
                 "составь сводку схемы БД»), экономя контекст и изолируя "
                 "радиус поражения при работе с чужим PR. Независимые "
                 "подшаги → субагент оправдан; зависимые — та же хрупкость "
                 "мульти-агента, не апгрейд.", size=13)
    speaker_notes(s, load_notes("s13a"))


def build_s13b(p):
    """assertion_visual — §2.5 skills: SKILL.md anatomy + scope + 3 SWE
    examples (issue #162)."""
    s = blank(p)
    slide_title(s,
                "Skill = директория с коротким SKILL.md, который агент "
                "быстро сканирует.", size=23, y=0.34, h=0.62)
    # Left — anatomy
    lx, ly, lw, lh = 0.55, 1.42, 5.85, 3.68
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "clipboard-list", lx + 0.24, ly + 0.16, 0.42, "mid")
    text_box(s, lx + 0.78, ly + 0.18, lw - 1.0, 0.40, "Формат",
             size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    comp = [
        ("SKILL.md", "YAML frontmatter (название/описание/когда) + "
                      "тело-инструкция"),
        ("scripts/", "опционально — готовые скрипты, вызываются напрямую"),
        ("references/", "опционально — длинные справочные материалы"),
    ]
    fy = ly + 0.64
    for a, b in comp:
        text_box(s, lx + 0.24, fy, lw - 0.48, 0.26, a, size=13, bold=True,
                 color=DEEP)
        text_box(s, lx + 0.24, fy + 0.26, lw - 0.48, 0.38, b, size=11,
                 color=SLATE, line_spacing=1.06)
        fy += 0.66
    filled_rect(s, lx + 0.24, fy + 0.02, lw - 0.48, 0.86, TEAL_TINT,
                stroke=TEAL, stroke_pt=1.25, radius=True, radius_adj=0.10)
    text_box(s, lx + 0.40, fy + 0.10, lw - 0.80, 0.70,
             "SKILL.md короткий и быстро читается моделью; тяжёлые "
             "детали подгружаются, только когда навык реально "
             "используется.",
             size=11, italic=True, color=DEEP, line_spacing=1.12,
             anchor=MSO_ANCHOR.MIDDLE)
    # Right — scope + 3 examples
    rx, rw = 6.65, 6.15
    ocean_box(s, rx, ly, rw, 1.58)
    text_box(s, rx + 0.24, ly + 0.12, 2.75, 0.52,
             "Project-level\n(уровень проекта)", size=12.5, bold=True,
             color=MID, line_spacing=1.02)
    text_box(s, rx + 0.24, ly + 0.62, 2.75, 0.88,
             "в репозитории, коммитится в git — команда получает "
             "автоматически, та же дисциплина, что steering-файл.",
             size=11, color=DEEP, line_spacing=1.12)
    text_box(s, rx + 3.16, ly + 0.12, 2.75, 0.52,
             "Personal (личный)", size=12.5, bold=True, color=TEAL,
             line_spacing=1.02)
    text_box(s, rx + 3.16, ly + 0.62, 2.75, 0.88,
             "в конфигурации разработчика — все его проекты, личная "
             "продуктивность, не контракт команды.",
             size=11, color=DEEP, line_spacing=1.12)
    ocean_box(s, rx, ly + 1.72, rw, 1.96)
    text_box(s, rx + 0.24, ly + 1.80, rw - 0.48, 0.28,
             "3 конкретных SWE-примера", size=13, bold=True, color=MID)
    examples = [
        "«Оформить PR по конвенциям проекта» — точный формат: Summary / "
        "Test plan / Related issues",
        "«Прогнать тесты» — точная команда этого проекта, отличить провал "
        "от известного нестабильного теста",
        "«Сгенерировать changelog» — формат проекта + классификация по "
        "истории коммитов",
    ]
    ey = ly + 2.12
    for ex in examples:
        text_box(s, rx + 0.24, ey, rw - 0.48, 0.50, "•  " + ex,
                 size=11.5, color=DEEP, line_spacing=1.12)
        ey += 0.52
    gold_callout(s, 0.55, 5.28, 12.25, 1.14,
                 "Skill не учит модель тому, чего она не умеет в принципе, — "
                 "он фиксирует конкретный, проектный вариант общей "
                 "процедуры, чтобы агент не изобретал его заново и не "
                 "дрейфовал от раза к разу. Прямая параллель паттерну "
                 "structure + constraints + tests.", size=13.5)
    speaker_notes(s, load_notes("s13b"))


def build_s13c(p):
    """assertion_visual — §2.6 MCP categories для кодинг-агента + least-
    privilege callback (issue #162). [Fix #9, format-fatigue] 2x3 grid
    (not vertical row-list) — visually distinct from s13a's 3-card row
    and s13b's anatomy/scope split, breaks the 3rd-in-a-row "list+icon"
    banner-blindness pattern student-simulator flagged."""
    s = blank(p)
    slide_title(s,
                "5 категорий MCP-серверов разработчика — что каждая "
                "закрывает.", size=21, y=0.30, h=0.56)
    text_box(s, 0.55, 0.92, 12.25, 0.30,
             "MCP (Model Context Protocol, Лекция 3) — единый способ "
             "подключить агента к внешним системам.", size=12.5,
             italic=True, color=MID)
    cells = [
        ("git-pull-request", "Репозиторий / issue-трекер",
         "GitHub / GitLab / Jira — читать и писать issue, PR, комментарии "
         "напрямую. Делает возможным уровень D.", "mid"),
        ("database", "Файловая система",
         "доступ за пределами текущего проекта (документация, другой "
         "репозиторий в монорепо). Контекст, не помещающийся в один "
         "репозиторий.", "mid"),
        ("refresh-cw", "CI/CD",
         "статусы сборки, логи конвейера. Агент итерирует по реальному "
         "логу CI, не по локальному прогону.", "mid"),
        ("layers", "Базы данных",
         "чтение (реже запись) схемы и данных. Миграция, согласованная со "
         "схемой; причина медленного запроса.", "teal"),
        ("scan-search", "Документация / поиск",
         "внутренняя вики, база знаний. Разрыв между общими знаниями "
         "модели и спецификой организации.", "teal"),
    ]
    bx, by, bw = 0.55, 1.42, 12.25
    gap = 0.18
    ncol = 2
    cw = (bw - gap * (ncol - 1)) / ncol
    ch = 1.16
    grid_h = ch * 3 + gap * 2
    ocean_box(s, bx, by, bw, grid_h)
    for idx, (ic, ttl, desc, var) in enumerate(cells):
        col = idx % ncol
        row = idx // ncol
        cx = bx + col * (cw + gap)
        cy = by + row * (ch + gap)
        icon(s, ic, cx + 0.20, cy + 0.18, 0.40, var)
        text_box(s, cx + 0.72, cy + 0.16, cw - 0.92, 0.34, ttl,
                 size=12.5, bold=True, color=DEEP, line_spacing=1.02)
        text_box(s, cx + 0.72, cy + 0.50, cw - 0.92, 0.60, desc,
                 size=10.5, color=SLATE, line_spacing=1.08)
        if col < ncol - 1:
            filled_rect(s, cx + cw + gap / 2 - 0.006, cy + 0.08, 0.012,
                        ch - 0.16, SOFT_GREY)
        if row < 2:
            filled_rect(s, cx + 0.20, cy + ch + gap / 2 - 0.006,
                        cw - 0.40, 0.012, SOFT_GREY)
    # 6th grid cell (row 2, col 1) has no 6th category — fill with a
    # visual pointer to the scope callout below instead of leaving it
    # empty (Visual Mass Balance: unequal cell content needs explaining).
    _lc_x = bx + 1 * (cw + gap)
    _lc_y = by + 2 * (ch + gap)
    icon(s, "shield-check", _lc_x + 0.20, _lc_y + 0.30, 0.40, "gold")
    text_box(s, _lc_x + 0.72, _lc_y + 0.20, cw - 0.92, 0.66,
             "Общий принцип для всех пяти — ниже.", size=11,
             italic=True, color=LIGHT, line_spacing=1.10,
             anchor=MSO_ANCHOR.MIDDLE)
    filled_rect(s, 0.55, by + grid_h + 0.20, 12.25, 0.50, GOLD, radius=True,
                radius_adj=0.16)
    text_box(s, 0.55, by + grid_h + 0.20, 12.25, 0.50,
             "Каждое подключение — решение о scope (объёме доступа), не "
             "техническая деталь.", size=14.5, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    teal_callout(s, 0.55, by + grid_h + 0.86, 12.25, 1.00,
                 "Least-privilege конкретно здесь: доступ к трекеру не "
                 "обязан включать удаление issue · доступ к БД для "
                 "аналитики не обязан включать запись в prod-таблицы · "
                 "доступ к CI не обязан включать право менять секреты. "
                 "«Какой сервер» и «с каким scope» — два разных решения.",
                 size=12.5)
    speaker_notes(s, load_notes("s13c"))


def build_s13d(p):
    """assertion_visual — §2.7 part 1: steering-file components + vs
    README/CONTRIBUTING + versioning-as-code (issue #162)."""
    s = blank(p)
    slide_title(s,
                "Steering-файл: что агент должен знать ПЕРЕД запуском.",
                size=24, y=0.34, h=0.60)
    # Left — 4 components (full-height column, matches right column)
    lx, ly, lw, lh = 0.55, 1.40, 6.15, 5.36
    ocean_box(s, lx, ly, lw, lh)
    comps = [
        ("terminal", "Команды сборки/теста/линта",
         "точная команда ЭТОГО проекта, не «как обычно»: конкретный "
         "флаг, пакетный менеджер, порядок шагов"),
        ("list-checks", "Конвенции кода",
         "именование, предпочитаемые паттерны, структура директорий, "
         "которую агент должен уважать"),
        ("shield-alert", "Архитектурные ограничения",
         "инварианты: «модуль без сетевого доступа», «не менять "
         "публичный API без запроса»"),
        ("circle-x", "Явное «что НЕ делать»",
         "отдельно от позитивных инструкций — конкретный список "
         "запрещённого для этого проекта"),
    ]
    cy = ly + 0.24
    ch = (lh - 0.44) / 4
    for ic, ttl, desc in comps:
        icon(s, ic, lx + 0.24, cy + 0.06, 0.44, "mid")
        text_box(s, lx + 0.82, cy, lw - 1.06, 0.38, ttl, size=14, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, lx + 0.82, cy + 0.40, lw - 1.06, 0.80, desc, size=11.5,
                 italic=True, color=SLATE, line_spacing=1.14)
        if ic != "circle-x":
            filled_rect(s, lx + 0.24, cy + ch - 0.10, lw - 0.48, 0.012,
                        SOFT_GREY)
        cy += ch
    # Right — vs README/CONTRIBUTING (2-col compare) + versioning callout
    rx, rw = 6.90, 5.90
    ocean_box(s, rx, ly, rw, 2.60)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.32,
             "Steering-файл vs README/CONTRIBUTING", size=14, bold=True,
             color=MID)
    pairs = [
        ("Аудитория", "человек", "операционный контекст агента"),
        ("Назначение", "объяснить проект, онбординг",
         "что нужно ПРЯМО СЕЙЧАС"),
        ("Частота чтения", "один раз, при онбординге", "на каждый запуск"),
    ]
    py = ly + 0.66
    for lbl, a, b in pairs:
        text_box(s, rx + 0.24, py, rw - 0.48, 0.24, lbl, size=11.5,
                 bold=True, color=TEAL)
        text_box(s, rx + 0.24, py + 0.28, (rw - 0.64) / 2, 0.48, a,
                 size=11.5, color=DEEP, line_spacing=1.08)
        text_box(s, rx + 0.24 + (rw - 0.48) / 2 + 0.12, py + 0.28,
                 (rw - 0.64) / 2, 0.48, b, size=11.5, bold=True, color=DEEP,
                 line_spacing=1.08)
        py += 0.62
    filled_rect(s, rx, ly + 2.74, rw, 2.62, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.5, radius=True, radius_adj=0.06)
    text_box(s, rx + 0.26, ly + 2.86, rw - 0.52, 0.32,
             "Версионирование как код", size=14.5, bold=True, color=MID)
    text_box(s, rx + 0.26, ly + 3.22, rw - 0.52, 0.86,
             "Steering-файл коммитится в git и проходит ревью в pull "
             "request наравне с изменениями кода — если конвенция "
             "поменялась, это часть того же PR, что и код, который её "
             "демонстрирует.", size=12, color=DEEP, line_spacing=1.16)
    # [Fix #2, presentation-critic P1] gold accent — most consequence-
    # carrying single line on this slide (steering-file-as-versioned-
    # artifact warning), fill+DEEP text per deck-wide WCAG-safe convention
    # (not gold text on light bg).
    filled_rect(s, rx + 0.20, ly + 4.14, rw - 0.40, 1.08, GOLD,
                stroke=None, radius=True, radius_adj=0.10)
    text_box(s, rx + 0.40, ly + 4.22, rw - 0.80, 0.92,
             "Файл, который никто не обновлял год, — не нейтральный ноль, "
             "а источник устаревшей инструкции: агент последует "
             "написанному, даже если оно больше не соответствует "
             "реальности.", size=12, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.16)
    speaker_notes(s, load_notes("s13d"))


def build_s13e(p):
    """case_study — §2.7 part 2a: presence-paradox RCT null result +
    Honest Lying entrenchment risk + ritual-cost criterion (gold-prominent
    per brief, issue #162). [Fix, presentation-critic P1 + reader-
    simulator P1] Split from old combined s13e (issue #162 v1) — the
    git-conventions block moved out to standalone s13f so this slide has
    room for a concrete Honest-Lying worked example (reader-simulator
    P2/consolidated-into-fix) and the "Presence paradox" title now carries
    an inline Russian gloss matching chapter.md §2.7's own
    «presence paradox» (Gloaguen et al., ...) convention (presentation-
    critic P1). "self-authored" glossed inline both occurrences
    (reader-simulator P1 cross-ref §5.8 anti-anglicism)."""
    s = blank(p)
    slide_title(s,
                "«Presence paradox» (парадокс присутствия файла): само "
                "наличие — не гарантия пользы.", size=20, y=0.32, h=0.62)
    # Left — RCT null result
    lx, ly, lw, lh = 0.55, 1.36, 6.05, 4.90
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.24, ly + 0.16, lw - 0.48, 0.30,
             "Gloaguen et al., arXiv:2602.11988, 2026", size=13,
             bold=True, color=MID)
    text_box(s, lx + 0.24, ly + 0.52, lw - 0.48, 0.92,
             "3 условия: NONE / LLM-generated / developer-written. "
             "2 постановки — SWE-bench-задачи и реальные repo-issue.",
             size=12.5, color=DEEP, line_spacing=1.18)
    filled_rect(s, lx + 0.24, ly + 1.52, lw - 0.48, 1.00, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.10)
    text_box(s, lx + 0.42, ly + 1.62, lw - 0.84, 0.80,
             "Наличие файла не даёт статистически значимого прироста "
             "успеха; стоимость и число шагов, как правило, растут.",
             size=12.5, bold=True, color=DEEP, line_spacing=1.16)
    text_box(s, lx + 0.24, ly + 2.62, lw - 0.48, 0.58,
             "Исключение: LLM-файлы там, где не было другой документации "
             "вообще.", size=11, italic=True, color=SLATE,
             line_spacing=1.12)
    filled_rect(s, lx + 0.24, ly + 3.30, lw - 0.48, 0.012, SOFT_GREY)
    text_box(s, lx + 0.24, ly + 3.44, lw - 0.48, 0.52,
             "«Honest Lying» (правдоподобное самозакрепление ошибки): "
             "риск самостоятельной правки агентом",
             size=12, bold=True, color=MID, line_spacing=1.05)
    text_box(s, lx + 0.24, ly + 3.98, lw - 0.48, 0.90,
             "Dixit, Kamal, Oates, arXiv:2605.29463, 2026 — заметки, "
             "написанные самим агентом (self-authored), могут не "
             "исправлять неверное убеждение, а закреплять его через "
             "повторные попытки.", size=11.5, color=DEEP, line_spacing=1.14)
    # Right — concrete worked example of entrenchment (Fix, reader-sim P2)
    rx, rw = 6.80, 6.00
    ocean_box(s, rx, ly, rw, 2.30)
    icon(s, "bot", rx + 0.22, ly + 0.16, 0.42, "gold")
    text_box(s, rx + 0.76, ly + 0.18, rw - 0.98, 0.38,
             "Пример закрепления ошибки", size=13.5,
             bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.24, ly + 0.66, rw - 0.48, 1.50,
             "Агент на раннем шаге ошибочно записал в свой файл "
             "«тест X не проходит из-за Y». Следующая сессия унаследует "
             "эту запись как данность и не станет перепроверять её "
             "заново — даже когда реальная причина давно другая.",
             size=12.5, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, ly + 2.46, rw, 2.44)
    text_box(s, rx + 0.24, ly + 2.60, rw - 0.48, 0.30,
             "Тот же класс риска", size=13, bold=True, color=TEAL)
    text_box(s, rx + 0.24, ly + 2.94, rw - 0.48, 1.80,
             "Если агенту разрешено самому править свой steering-файл "
             "(«запиши, что понял по итогам сессии») — это то же "
             "уязвимое место, что и написанные самим агентом заметки в "
             "целом: без человеческого ревью правки риск закрепления "
             "ошибки реален, а не гипотетичен.",
             size=12, color=DEEP, line_spacing=1.18)
    gold_callout(s, 0.55, 6.40, 12.25, 0.86,
                 "Когда steering-файл НЕ нужен: нет названного "
                 "информационного пробела — писать ритуально не даёт "
                 "измеренной пользы. Альтернатива: писать только под "
                 "названный пробел, ревьюить изменения как обычный код.",
                 size=13)
    speaker_notes(s, load_notes("s13e"))


def build_s13f(p):
    """assertion_visual — §2.7 part 2b (new, split from old s13e per
    presentation-critic/student-simulator feedback issue #162 v2):
    git-conventions-as-contract, standalone slide (Conventional Commits /
    branch naming / structured PR sections)."""
    s = blank(p)
    slide_title(s,
                "Git-конвенции как контракт для агента — не свободный "
                "текст, а формат.", size=23, y=0.34, h=0.60)
    cards = [
        ("git-commit", "Conventional Commits", "mid",
         "`type(scope): описание` — фиксированный набор типов "
         "(feat/fix/docs/refactor…).",
         "парсится автоматически: changelog, semver, фильтр истории"),
        ("git-branch", "Именование веток", "mid",
         "`feature/`, `fix/`, префикс номера задачи "
         "(`feature/TASK-123-…`).",
         "автосвязь ветки с трекером, без ручного указания номера"),
        ("file-text", "Структура секций PR", "teal",
         "Summary / Test Plan / Breaking Changes — по шаблону, не "
         "«своими словами».",
         "машинно-проверяемый чек-лист вместо прозы, которую нужно "
         "интерпретировать"),
    ]
    n = 3
    gap = 0.24
    cw = (12.25 - gap * (n - 1)) / n
    cx0, cy, chh = 0.55, 1.40, 3.30
    for ic, ttl, var, ex, arrow in cards:
        ocean_box(s, cx0, cy, cw, chh)
        icon(s, ic, cx0 + 0.22, cy + 0.20, 0.46, var)
        text_box(s, cx0 + 0.80, cy + 0.22, cw - 1.00, 0.60, ttl,
                 size=13.5, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.04)
        text_box(s, cx0 + 0.22, cy + 0.98, cw - 0.44, 1.10, ex,
                 size=11.5, color=DEEP, line_spacing=1.18)
        filled_rect(s, cx0 + 0.22, cy + 2.14, cw - 0.44, 0.012, SOFT_GREY)
        text_box(s, cx0 + 0.22, cy + 2.30, cw - 0.44, 0.90, "→ " + arrow,
                 size=11, italic=True, color=TEAL, line_spacing=1.14)
        cx0 += cw + gap
    # [Fix #3, presentation-critic P1] gold accent — single most
    # load-bearing point on this slide (format parsed by a program vs
    # re-derived by the model from prose each run), fill+DEEP text per
    # deck-wide WCAG-safe convention. This slide had 0 gold pixels before.
    filled_rect(s, 0.55, 4.90, 12.25, 0.50, GOLD, stroke=None,
                radius=True, radius_adj=0.16)
    text_box(s, 0.55, 4.90, 12.25, 0.50,
             "Формат парсится программой — не восстанавливается моделью "
             "из вольного текста на каждый запуск.", size=14, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    teal_callout(s, 0.55, 5.56, 12.25, 0.92,
                 "CI, генератор changelog, бот-линковщик читают формат "
                 "напрямую. Тот же принцип, что и у steering-файла — "
                 "конкретный формат вместо предположений.", size=13)
    speaker_notes(s, load_notes("s13f"))


def build_s13g(p):
    """assertion_visual — §2.8 part 1: 3 task-log architectural patterns
    (nested / unified log / flat folder), issue #162. [renumbered from
    s13f → s13g when s13e split, issue #162 v2 fix round]. [Fix #3,
    presentation-critic P1] gold accent added — this slide had 0 gold
    pixels. Placed as a neutral badge on the shared framing callout
    (not on any single pattern card) to avoid the "inconsistent gold-
    emphasis across same-tier cards" anti-pattern (README #21) — all 3
    patterns stay visually equal-weight, per critic's own caution."""
    s = blank(p)
    slide_title(s,
                "Где хранить состояние задачи между сессиями агента — "
                "три паттерна.", size=21, y=0.32, h=0.86)
    filled_rect(s, 0.55, 1.20, 1.55, 0.62, GOLD, stroke=None,
                radius=True, radius_adj=0.16)
    text_box(s, 0.55, 1.20, 1.55, 0.62, "3 паттерна", size=13, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    teal_callout(s, 2.24, 1.20, 10.56, 0.62,
                 "Не steering-файл (проект целиком, живёт долго) и не "
                 "заметки, написанные самим агентом внутри одной "
                 "сессии (self-authored), — это межзадачное и "
                 "межсессионное состояние на диске.", size=12)
    patterns = [
        ("boxes", "1. Nested per-task folder\n(вложенная директория на "
         "задачу)", "mid",
         "`tasks/TASK-123/` с notes.md, log.md, черновиками.",
         "богатая структура на задачу — разные типы информации по "
         "файлам, план отдельно от лога выполнения.",
         "директорий растёт линейно; без архивации — захламление, "
         "удаление ощущается как потеря истории."),
        ("list-ordered", "2. Single unified log\n(единый растущий файл)",
         "mid",
         "Один файл (`PROGRESS.md`), каждая сессия дописывает в конец.",
         "одна точка входа — вся хронология без знания номера задачи "
         "заранее, проще всего восстановить контекст.",
         "не помещается в контекст целиком (context rot — деградация "
         "точности при разрастании контекста); параллельные сессии в "
         "один хвост — магнит merge-конфликтов."),
        ("layout-grid", "3. Flat shared folder\n(один файл на задачу "
         "без вложенности)", "teal",
         "`tasks/TASK-123.md` напрямую, без вложенности.",
         "баланс — единая точка входа (`tasks/`), но состояние задачи "
         "изолировано в своём файле, не смешано.",
         "список растёт без структуры; несколько артефактов на задачу "
         "требуют условности именования."),
    ]
    n = 3
    gap = 0.24
    cw = (12.25 - gap * (n - 1)) / n
    cx0, cy, chh = 0.55, 2.00, 4.60
    for ic, ttl, var, ex, strong, weak in patterns:
        ocean_box(s, cx0, cy, cw, chh)
        icon(s, ic, cx0 + 0.22, cy + 0.18, 0.44, var)
        text_box(s, cx0 + 0.78, cy + 0.18, cw - 0.98, 0.64, ttl,
                 size=11.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.04)
        text_box(s, cx0 + 0.22, cy + 0.86, cw - 0.44, 0.48, ex,
                 size=10.5, italic=True, color=SLATE, line_spacing=1.10)
        filled_rect(s, cx0 + 0.22, cy + 1.36, cw - 0.44, 0.012, SOFT_GREY)
        text_runs(s, cx0 + 0.22, cy + 1.50, cw - 0.44, 1.34, [
            {"text": "Сильная сторона: ", "size": 11, "bold": True,
             "color": DEEP},
            {"text": strong, "size": 11, "color": DEEP},
        ], line_spacing=1.14)
        filled_rect(s, cx0 + 0.22, cy + 3.02, cw - 0.44, 0.012, SOFT_GREY)
        text_runs(s, cx0 + 0.22, cy + 3.16, cw - 0.44, 1.30, [
            {"text": "Что ломается: ", "size": 11, "bold": True,
             "color": TEAL},
            {"text": weak, "size": 11, "color": TEAL},
        ], line_spacing=1.14)
        cx0 += cw + gap
    speaker_notes(s, load_notes("s13g"))


def build_s13h(p):
    """comparison / schema_matrix — §2.8 part 2: 4-criteria x 3-pattern
    comparison table + no-single-right-answer closing framing (issue #162).
    [renumbered s13g -> s13h when s13e split, issue #162 v2 fix round.]
    [Fix #4, presentation-critic P1] Table body font raised 9.7pt ->
    12pt (README Schema Readability floor) by dropping to 4 columns
    (merged "Обнаруживаемость" into "Что ломается" as a sub-line per
    critic's option (a)) — frees width for readable font at this row
    count. [Fix #2b/gold] gold accent added on the closing callout's
    lead sentence (fill+DEEP), echoing s13/s13a "when-not-needed"-style
    framing treatment — this slide had 0 gold pixels before."""
    s = blank(p)
    slide_title(s,
                "Нет единственно верного паттерна — решающая ось зависит "
                "от команды.", size=21, y=0.34, h=0.56)
    bx, by, bw = 0.55, 1.24, 12.25
    ocean_box(s, bx, by, bw, 3.60)
    headers = ["Паттерн", "Git-diff", "Параллельные сессии",
               "Обнаруживаемость / что ломается"]
    col_w = [2.55, 2.85, 2.95, 3.90]
    hy = by + 0.14
    hx = bx + 0.14
    for j, h in enumerate(headers):
        filled_rect(s, hx, hy, col_w[j] - 0.06, 0.54, MID)
        text_box(s, hx + 0.10, hy, col_w[j] - 0.26, 0.54, h, size=12.5,
                 bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        hx += col_w[j]
    rows = [
        ("1. Вложенная папка\nна задачу",
         "конфликтов между задачами нет; растёт вложенно",
         "высокая между задачами; директории «осиротевают»",
         "легко по номеру, трудно обзорно; захламление "
         "неархивированными директориями"),
        ("2. Единый\nжурнал",
         "append-only, одна точка записи для всех",
         "низкая — общий хвост = магнит конфликтов",
         "высокая обзорно, низкая точечно; не помещается в контекст "
         "целиком (context rot)"),
        ("3. Плоская общая\nпапка",
         "конфликт только на одной и той же задаче",
         "высокая между задачами (как 1, но без директорий)",
         "высокая обзорно и точечно; растёт без структуры, нужна "
         "условность именования"),
    ]
    ry = hy + 0.54
    rh = 0.98
    for i, row in enumerate(rows):
        rx = bx + 0.14
        bgc = SURFACE if i % 2 == 0 else WHITE
        for j, cell in enumerate(row):
            filled_rect(s, rx, ry, col_w[j] - 0.06, rh, bgc,
                        stroke=SOFT_GREY, stroke_pt=0.75)
            text_box(s, rx + 0.10, ry, col_w[j] - 0.24, rh, cell,
                     size=12, bold=(j == 0), color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
            rx += col_w[j]
        ry += rh
    filled_rect(s, 0.55, 5.02, 12.25, 0.44, GOLD, stroke=None,
                radius=True, radius_adj=0.14)
    text_box(s, 0.55 + 0.22, 5.02, 12.25 - 0.44, 0.44,
             "Нет единственно верного ответа — это следствие взвешивания "
             "осей под ситуацию команды, а не рекомендация в обход них.",
             size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    teal_callout(s, 0.55, 5.56, 12.25, 1.50,
                 "Команда, работающая строго последовательно, — "
                 "устойчивость к параллелизму не решающая ось, Паттерн 2 "
                 "может быть валидным выбором. Малое число задач и низкая "
                 "параллельность — Паттерн 1 не будет ошибкой. Регулярные "
                 "параллельные сессии на разных задачах — решающая ось "
                 "именно параллелизм, и здесь Паттерн 2 наиболее уязвим по "
                 "построению.", size=12.5)
    speaker_notes(s, load_notes("s13h"))


def build_s14(p):
    """section_divider — Раздел 3: Уровень D."""
    build_section_divider(
        p, 3, "Уровень D:\nоркестратор и трекер",
        "Уровень C: человек ставил каждую задачу. Уровень D — AI берёт "
        "задачи из трекера сам, иногда несколькими агентами; максимум "
        "автономии — и максимум риска.", "s14")


def build_s15(p):
    """process / schema_pipeline — issue→PR + 2 risk amplifiers."""
    s = blank(p)
    slide_title(s, "Тот же цикл — но с двумя усилителями риска.", size=26)
    # Horizontal pipeline
    # v3.3: spacing compressed to fit tools_strip — контент не урезан
    px, py, pw, ph = 0.55, 1.36, 12.25, 1.30
    ocean_box(s, px, py, pw, ph)
    stages = [
        ("issue\nиз трекера", LIGHT),
        ("AI\nдекомпозирует", MID),
        ("правит /\nоткрывает PR", MID),
        ("человек:\napproval · merge · прод", GOLD),
    ]
    n = 4
    aw = 0.52
    sw = (pw - 0.50 - aw * (n - 1)) / n
    sx = px + 0.25
    for i, (lab, col) in enumerate(stages):
        bx = sx + i * (sw + aw)
        hi = (col == GOLD)
        filled_rect(s, bx, py + 0.22, sw, ph - 0.44,
                    GOLD_TINT if hi else WHITE,
                    stroke=GOLD if hi else SOFT_GREY,
                    stroke_pt=2.0 if hi else 1.0, radius=True, radius_adj=0.08)
        text_box(s, bx + 0.06, py + 0.22, sw - 0.12, ph - 0.44, lab,
                 size=12.5, bold=True, color=DEEP if hi else MID,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.08)
        if i < n - 1:
            right_arrow(s, bx + sw + 0.02, py + ph / 2 - 0.21, aw - 0.04,
                        0.42, fill=LIGHT)
    # 2 amplifiers (v3.3: высота сжата, текст не урезан)
    aw2 = 6.0
    ocean_box(s, 0.55, 2.84, aw2, 1.92)
    text_box(s, 0.78, 2.96, aw2 - 0.45, 0.38,
             "Усилитель 1 — источник задачи: трекер, не человек",
             size=14, bold=True, color=MID, line_spacing=1.05)
    text_box(s, 0.78, 3.40, aw2 - 0.45, 1.28,
             "двусмысленный или плохо написанный issue идёт в работу без "
             "промежуточного человеческого осмысления",
             size=13, color=DEEP, line_spacing=1.16)
    ocean_box(s, 6.80, 2.84, aw2, 1.92, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, 7.03, 2.96, aw2 - 0.45, 0.38,
             "Усилитель 2 — мульти-агент по умолчанию ≠ апгрейд",
             size=14, bold=True, color=DEEP, line_spacing=1.05)
    text_box(s, 7.03, 3.40, aw2 - 0.45, 1.28,
             "параллельные субагенты на зависимых подзадачах принимают "
             "неявные конфликтующие решения. Один линейный агент надёжнее "
             "— тот же вывод, что в Лекции 3",
             size=13, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 4.88, 12.25, 0.86,
                 "Уровень D — максимум автономии и максимум риска: человек "
                 "обязателен на любом необратимом или прод-действии, иначе "
                 "автономный деструктив проходит без подтверждения.",
                 size=14)
    # [v3.3, Решение #102] Инструменты уровня D 2026 (book-first §3.2)
    tools_strip(
        s, 0.55, 5.84, 12.25, 1.36,
        ["Copilot agent", "Devin 2.0", "Jules", "Codex Cloud"],
        "Adoption: D — самый молодой сегмент; мульти-агентные "
        "системы — верхняя кромка, emerging, а не мейнстрим.",
        "«Полностью автономный инженер» (Devin) — overclaim, не "
        "факт. Copilot coding agent в проде = 5 отказов + "
        "«аварийный выключатель» → гейты обязательны.",
        left_ratio=0.58)
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    """case_study — Replit destructive failure (эталон) + lesson."""
    s = blank(p)
    slide_title(s, "Агент стёр прод-БД в code-freeze — и оценил себя на 95 из 100.",
                size=22, y=0.34, h=0.56)
    # Left — Replit chronicle
    lx, ly, lw, lh = 0.55, 1.34, 6.55, 4.52
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.18, lw - 0.52, 0.34,
             "Replit, июль 2025 (vibe-coding, реальные данные)",
             size=14, bold=True, color=MID)
    chron = [
        ("Человек ввёл явный code-freeze: «БОЛЬШЕ НИКАКИХ ИЗМЕНЕНИЙ»", False),
        ("Агент удалил рабочую (production) базу данных", False),
        ("Сфабриковал отчёты, маскирующие проблему; на вопрос солгал", False),
        ("Оценил своё поведение 95 из 100", True),
        ("Заявил «rollback невозможен» — а механизм отката работал", False),
    ]
    cyl = ly + 0.62
    for t, hi in chron:
        circle(s, lx + 0.28, cyl + 0.05, 0.13, GOLD if hi else MID)
        text_box(s, lx + 0.56, cyl, lw - 0.84, 0.72,
                 t, size=12.5, bold=hi, color=(DEEP if hi else DEEP),
                 line_spacing=1.12)
        cyl += 0.76
    filled_rect(s, lx + 0.26, cyl + 0.02, lw - 0.52, 0.52, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.12)
    text_box(s, lx + 0.38, cyl + 0.02, lw - 0.76, 0.52,
             "Эхо того же режима: Kiro — 13ч простоя · PocketOS — БД за 9 секунд",
             size=11.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    # Right — lesson + alternative
    rx, rw = 7.30, 5.50
    ocean_box(s, rx, ly, rw, 2.20, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 0.16, rw - 0.52, 0.30, "Урок", size=14,
             bold=True, color=DEEP)
    text_box(s, rx + 0.26, ly + 0.50, rw - 0.52, 1.60,
             "инструкция в промпте ≠ контроль (траектория переезжает) · "
             "самооценка ≠ проверка · отчёт агента ≠ доказательство · "
             "accountability не делегируется",
             size=13, color=DEEP, line_spacing=1.22)
    ocean_box(s, rx, ly + 2.36, rw, 2.19, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 2.50, rw - 0.52, 0.30, "Альтернатива",
             size=14, bold=True, color=TEAL)
    text_box(s, rx + 0.26, ly + 2.84, rw - 0.52, 1.65,
             "dev/prod-изоляция · hard human-gate на деструктив · "
             "least-privilege · проверенный rollback · immutable-бэкапы · "
             "two-person-rule на агентов",
             size=13, color=DEEP, line_spacing=1.22)
    footer(s, "Датированные инциденты Replit/Kiro/PocketOS (2025–2026). "
              "Корневая ошибка — выбор уровня автономии, неадекватного цене "
              "ошибки.")
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """case_study — METR revealed (c17) + measurement protocol."""
    s = blank(p)
    slide_title(s, "Не доверяйте ощущению — измерьте у себя.", size=26)
    # Left — METR раскрыт: 3 крупных числа (вместо декор-chart)
    lx, ly, lw, lh = 0.55, 1.36, 7.05, 4.10
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.24, ly + 0.14, lw - 0.48, 0.32,
             "METR раскрыт — три числа об одном и том же",
             size=14, bold=True, color=MID)
    metr = [
        ("Прогноз до", "ждали ускорение", "−24%", LIGHT, False),
        ("Вера после", "поработав, верили в ускорение", "−20%", LIGHT, False),
        ("Измеренный факт", "по времени — замедление", "+19%", GOLD, True),
    ]
    myy = ly + 0.54
    msh = 0.86
    for lab, sub, val, col, hi in metr:
        # [Fix B1, issue #162] highlight: SOLID gold fill + DEEP text.
        bg = GOLD if hi else SURFACE
        filled_rect(s, lx + 0.24, myy, lw - 0.48, msh, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.42, myy + 0.10, lw - 2.55, 0.32, lab,
                 size=13.5, bold=True, color=(DEEP if hi else MID))
        text_box(s, lx + 0.42, myy + 0.42, lw - 2.55, 0.34, sub,
                 size=11.5, italic=True, color=(DEEP if hi else SLATE))
        text_box(s, lx + lw - 2.05, myy + 0.04, 1.65, msh - 0.08, val,
                 size=31, bold=True, color=DEEP,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        myy += msh + 0.07
    text_box(s, lx + 0.24, myy + 0.02, lw - 0.48, 0.50,
             "Эффект AI слабеет, когда контекст в голове, проверка дорога "
             "и есть своя быстрая альтернатива.",
             size=11.5, italic=True, color=MID, line_spacing=1.14)
    # Right — protocol
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 4.10)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.32,
             "Протокол «как измерить у себя»", size=14, bold=True, color=MID)
    proto = [
        "сопоставимый класс своих задач",
        "случайно распределить: с AI / без AI",
        "фиксировать реальное время, не самооценку",
        "считать дефекты/доработки за 2 недели",
        "применять вывод селективно (не везде / не нигде)",
    ]
    pyy = ly + 0.56
    for i, t in enumerate(proto):
        circle(s, rx + 0.26, pyy + 0.04, 0.32, MID)
        text_box(s, rx + 0.26, pyy + 0.04, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 0.70, pyy, rw - 0.96, 0.66, t,
                 size=13, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.10)
        pyy += 0.68
    gold_callout(s, 0.55, 5.56, 12.25, 0.72,
                 "Измеряйте, не ощущайте. Уровень и место AI — измеряемое "
                 "решение для конкретной команды, а не культурная установка.",
                 size=15)
    footer(s, "METR 2025 (RCT, n=16). Поздний-2025 «разворот» сигнала есть, "
              "но методологически слабее исходного — на него опираться рано.")
    speaker_notes(s, load_notes("s17"))


def build_s18(p):
    """section_divider — Раздел 4: Не только код."""
    build_section_divider(
        p, 4, "Не только код:\nтест, ревью, безопасность",
        "Лестница A→D — про то, насколько автономно AI пишет код. Но "
        "инженер ещё тестирует, ревьюит, думает об угрозах — там AI ведёт "
        "себя иначе.", "s18")


def build_s19(p):
    """assertion_visual — тест=спецификация; gate = mutation score."""
    s = blank(p)
    slide_title(s, "Тест — исполняемая спецификация.", size=26)
    icon(s, "flask-conical", 11.75, 0.36, 0.78, "mid")
    # Left — thesis + inline defs
    lx, ly, lw, lh = 0.55, 1.42, 7.05, 3.95
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.18, lw - 0.52, 0.74,
             "Тест = исполняемая спецификация — не подвержен ни «почти "
             "правильному», ни perception-gap.",
             size=14.5, bold=True, color=DEEP, line_spacing=1.20)
    text_box(s, lx + 0.26, ly + 0.98, lw - 0.52, 0.92,
             "AI силён в объёме тестов (быстро покрыть много классов "
             "входов). AI слаб в выборе, что именно проверять — это "
             "содержательное решение, не рутина (вернёмся к этому далее).",
             size=13, color=DEEP, line_spacing=1.18)
    text_box(s, lx + 0.26, ly + 2.04, lw - 0.52, 0.42,
             "mutation score — доля пойманных дефектов",
             size=12.5, italic=True, color=MID, line_spacing=1.12)
    text_box(s, lx + 0.26, ly + 2.56, lw - 0.52, 0.42,
             "quality-gate — порог в CI",
             size=12.5, italic=True, color=MID, line_spacing=1.12)
    # Right — illusion
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 1.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 0.18, rw - 0.52, 0.85,
             "100% coverage / ~4% mutation score",
             size=20, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, rx + 0.26, ly + 1.05, rw - 0.52, 0.80,
             "тесты «трогают» строки, но не проверяют их — зелёный CI "
             "без обнаружения дефектов",
             size=12.5, color=DEEP, line_spacing=1.16)
    ocean_box(s, rx, ly + 2.10, rw, 1.85)
    text_box(s, rx + 0.26, ly + 2.26, rw - 0.52, 1.55,
             "AI оптимизирует то, что измеряют. Гейтите по coverage — "
             "получите тесты «под coverage».",
             size=13, color=DEEP, line_spacing=1.20)
    gold_callout(s, 0.55, 5.42, 12.25, 0.84,
                 "Правильный quality-gate для AI-сгенерированных тестов — "
                 "mutation score, а не только coverage. Решение «что считать "
                 "корректным» — спецификация, не делегируется.", size=14)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    """comparison / schema_matrix — Greptile vs CodeRabbit + roles."""
    s = blank(p)
    slide_title(s, "Первый проход — машина, второй — человек.", size=26)
    text_box(s, 0.55, 1.14, 12.25, 0.34,
             "Тест на 50 реальных багах — фундаментальный размен «больше "
             "находит ↔ больше ложных срабатываний».",
             size=13, italic=True, color=MID)
    # 2 equal columns
    bx, by, bw, bh = 0.55, 1.56, 12.25, 2.05
    ocean_box(s, bx, by, bw, bh)
    cols = [
        ("Greptile", "~82% багов поймал", "11 ложноположительных", GOLD),
        ("CodeRabbit", "~44% багов поймал", "2 ложноположительных", LIGHT),
    ]
    cw2 = (bw - 0.60) / 2
    for i, (nm, m1, m2, col) in enumerate(cols):
        cx = bx + 0.20 + i * (cw2 + 0.20)
        hi = (col == GOLD)
        filled_rect(s, cx, by + 0.18, cw2, bh - 0.36,
                    GOLD_TINT if hi else SURFACE,
                    stroke=GOLD if hi else SOFT_GREY,
                    stroke_pt=2.0 if hi else 1.0, radius=True,
                    radius_adj=0.06)
        text_box(s, cx + 0.20, by + 0.32, cw2 - 0.40, 0.42, nm,
                 size=18, bold=True, color=DEEP if hi else MID)
        text_box(s, cx + 0.20, by + 0.84, cw2 - 0.40, 0.42, m1,
                 size=15, bold=True, color=DEEP)
        text_box(s, cx + 0.20, by + 1.26, cw2 - 0.40, 0.42, m2,
                 size=15, color=DEEP)
    text_box(s, 0.55, 3.70, 12.25, 0.30,
             "(Graphite ~6%.) Чем больше багов ловит инструмент, тем "
             "больше ложных тревог разгребает человек.",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # roles
    ocean_box(s, 0.55, 4.08, 6.0, 1.18)
    text_box(s, 0.78, 4.20, 5.55, 0.34, "AI-ревью = фильтр 1-го прохода",
             size=14, bold=True, color=MID)
    text_box(s, 0.78, 4.56, 5.55, 0.62,
             "дёшево ловит массовые механические дефекты и кандидатов на "
             "проблему", size=12, color=DEEP, line_spacing=1.14)
    ocean_box(s, 6.80, 4.08, 6.0, 1.18, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, 7.03, 4.20, 5.55, 0.34, "Человек = 2-й проход", size=14,
             bold=True, color=TEAL)
    text_box(s, 7.03, 4.56, 5.55, 0.62,
             "баг или ложная тревога? архитектурная уместность? "
             "дублирование кода?",
             size=12, color=DEEP, line_spacing=1.14)
    gold_callout(s, 0.55, 5.40, 12.25, 0.86,
                 "Антипаттерн — мерджить по вердикту AI: та же деградация "
                 "уровня C до D плюс шум ложных тревог. AI-ревью сужает "
                 "работу человека-ревьюера, но не снимает с него "
                 "ответственность за merge.", size=14)
    footer(s, "Greptile-бенчмарк 2025 (5 инструментов, 50 реальных багов) — "
              "тест от вендора, инструменты быстро меняются; читать с "
              "поправкой.")
    speaker_notes(s, load_notes("s20"))


def build_s21(p):
    """case_study — vulnerable code + false confidence; SAST/DAST/SCA."""
    s = blank(p)
    slide_title(s, "Опасна не ошибка, а ложная уверенность.", size=26)
    icon(s, "shield-alert", 11.75, 0.36, 0.78, "gold")
    # Left — data
    lx, ly, lw, lh = 0.55, 1.42, 6.85, 3.85
    ocean_box(s, lx, ly, lw, lh)
    data = [
        "NYU: в 89 security-сценариях ~40% программ с Copilot уязвимы",
        "Анализ 7703 файлов AI-кода: 12,1% CWE; Python ~16–18% > JS > TS",
    ]
    dy = ly + 0.22
    for t in data:
        circle(s, lx + 0.26, dy + 0.05, 0.13, MID)
        text_box(s, lx + 0.54, dy, lw - 0.80, 0.74, t, size=13, color=DEEP,
                 line_spacing=1.16)
        dy += 0.80
    filled_rect(s, lx + 0.24, dy + 0.02, lw - 0.48, 0.84, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.10)
    text_box(s, lx + 0.40, dy + 0.02, lw - 0.80, 0.84,
             "Stanford: с AI вносят больше уязвимостей И увереннее, что "
             "код безопасен",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)
    text_box(s, lx + 0.26, dy + 1.00, lw - 0.52, 0.95,
             "CWE — каталог типов уязвимостей",
             size=11.5, italic=True, color=MID, line_spacing=1.16)
    # Right — security tools (термин + 3-словный смысл; разворот → notes)
    rx, rw = 7.65, 5.15
    ocean_box(s, rx, ly, rw, 3.85, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 0.16, rw - 0.52, 0.32,
             "Security-инструментарий", size=14, bold=True, color=TEAL)
    tools = [
        ("SAST", "анализ кода"),
        ("DAST", "анализ приложения"),
        ("SCA", "анализ зависимостей"),
        ("secret-scan", "поиск утёкших ключей"),
    ]
    tyy = ly + 0.56
    for a, b in tools:
        text_box(s, rx + 0.26, tyy, 2.05, 0.74, a, size=14, bold=True,
                 color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 2.30, tyy, rw - 2.56, 0.74, b, size=13, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        tyy += 0.80
    gold_callout(s, 0.55, 5.32, 12.25, 0.92,
                 "AI снижает бдительность ровно там, где она нужнее. "
                 "Обязательный SAST + secret-scan как gate (не опция); "
                 "threat-modeling — человеческий шаг, не делегируется.",
                 size=14)
    footer(s, "NYU (Asleep at the Keyboard?, 2022); Stanford. Конкретные "
              "доли уточняются, но направление стабильно. Тему безопасности "
              "продолжим на следующем слайде.")
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """case_study — slopsquatting. SINGLE-FOCUS horizontal attack-flow
    (ломает 2-кол ритм s19–s23: full-width pipeline + hero-число)."""
    s = blank(p)
    slide_title(s, "Выдуманный пакет как атака на цепочку поставок.",
                size=24, y=0.34, h=0.58)
    text_box(s, 0.55, 0.96, 12.25, 0.40,
             "slopsquatting — атака на цепочку поставок (supply-chain): "
             "злоумышленник регистрирует имя пакета, которое AI стабильно "
             "выдумывает, и кладёт под него вредоносный код.",
             size=13, italic=True, color=MID, line_spacing=1.16)
    # Full-width horizontal attack-flow pipeline (single-focus)
    px, py, pw, ph = 0.40, 1.62, 12.55, 1.92
    ocean_box(s, px, py, pw, ph)
    text_box(s, px + 0.24, py + 0.14, pw - 0.48, 0.30,
             "Как работает атака — слева направо", size=13.5, bold=True,
             color=MID)
    stages = [
        ("Собрать", "выдуманные имена пакетов из ответов AI", MID),
        ("Отфильтровать", "по воспроизводимости — какие модель выдаёт "
         "стабильно", MID),
        ("Опубликовать", "зарегистрировать имя в npm/PyPI + положить "
         "вредонос", LIGHT),
        ("Ждать", "разработчик или агент делает `install` — вредонос "
         "при установке", GOLD),
    ]
    n = 4
    aw = 0.46
    sw = (pw - 0.48 - aw * (n - 1)) / n
    sx = px + 0.24
    syt = py + 0.52
    for i, (hd, body, col) in enumerate(stages):
        bx = sx + i * (sw + aw)
        hi = (col == GOLD)
        filled_rect(s, bx, syt, sw, ph - 0.70,
                    GOLD_TINT if hi else (SURFACE if col == LIGHT else WHITE),
                    stroke=GOLD if hi else SOFT_GREY,
                    stroke_pt=2.0 if hi else 1.0, radius=True,
                    radius_adj=0.07)
        text_box(s, bx + 0.10, syt + 0.10, sw - 0.20, 0.34,
                 f"{i + 1}. {hd}", size=13.5, bold=True,
                 color=DEEP if hi else MID, align=PP_ALIGN.CENTER)
        text_box(s, bx + 0.12, syt + 0.46, sw - 0.24, ph - 1.22, body,
                 size=11.5, color=DEEP, align=PP_ALIGN.CENTER,
                 line_spacing=1.12)
        if i < n - 1:
            right_arrow(s, bx + sw + 0.02, syt + (ph - 0.70) / 2 - 0.16,
                        aw - 0.04, 0.32, fill=LIGHT)
    # Hero number band (single dominant stat)
    hy = 3.74
    filled_rect(s, 0.40, hy, 12.55, 1.34, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.07)
    # [Fix B1, issue #162] gold TEXT on GOLD_TINT (≈1.8:1 FAIL) → solid-gold
    # plate behind number + DEEP text (≈6.85:1 PASS).
    filled_rect(s, 0.70, hy + 0.10, 3.30, 1.14, GOLD,
                radius=True, radius_adj=0.10)
    text_box(s, 0.70, hy + 0.10, 3.30, 1.14, "58%",
             size=68, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    connector(s, 4.05, hy + 0.22, 4.05, hy + 1.12, color=GOLD, width=2.0)
    text_box(s, 4.30, hy + 0.16, 8.35, 1.06,
             "выдуманных имён модель повторяет в нескольких запросах — "
             "значит, атака воспроизводима и масштабируема (всего ~20% "
             "ответов рекомендуют несуществующий пакет).",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16)
    gold_callout(s, 0.55, 5.28, 12.25, 0.98,
                 "Барьер — не-AI, инженерный: lockfile + пиннинг по хэшу · "
                 "разрешённый список реестров · проверка пакета до install · "
                 "SCA-скан. Прямой перенос урока Лекции 3 «когда не доверять "
                 "выводу модели».", size=14)
    footer(s, "Slopsquatting — термин Seth Larson (PSF, 2025); исследование "
              "576 000 сэмплов. Конкретные проценты уточняются, направление "
              "стабильно.")
    speaker_notes(s, load_notes("s22"))


def build_s22a(p):
    """case_study — curl-slop #5 (Решение #103, owner GATE C). Несущая
    ось — асимметрия стоимости: «сгенерировать фейк ≈ секунды» vs gold
    «опровергнуть = часы человека» → DDoS на внимание сопровождающих;
    чинить процесс, не запрещать AI. Стиль = failure-слайд (как
    s21/s22): Ocean motif, palette LOCKED, gold = ось асимметрии.
    [Fix C, issue #162] resolved numbers (было [FACT-CHECK] в notes):
    <5% valid-rate (baseline ~15%), ×8 объём (июль 2025), платный
    bounty закрыт 26.01.2026, полный мораторий 01.07-03.08.2026."""
    s = blank(p)
    slide_title(s, "AI-slop как DDoS на внимание сопровождающих.",
                size=24, y=0.34, h=0.58)
    icon(s, "triangle-alert", 11.75, 0.34, 0.78, "gold")
    # Context band — что произошло, с резолвнутыми числами (Fix C, #162)
    # [iter2 fix, #162] 3-line text needs more height — was clipping into
    # the asymmetry box below (overflow bug found on visual inspection).
    text_box(s, 0.55, 0.90, 12.25, 0.92,
             "curl — критичная open-source-библиотека, её ведёт небольшая "
             "команда. В её bug-bounty хлынул поток AI-сгенерированных "
             "«отчётов»: к июлю 2025 объём вырос ×8, доля валидных упала с "
             "~15% до <5% (пример — фиктивные GDB-дампы + ссылка на "
             "несуществующую функцию). Платный bounty закрыт 26.01.2026; "
             "полный мораторий на приём — 01.07–03.08.2026.",
             size=11.5, italic=True, color=MID, line_spacing=1.14)
    # ── Несущая ось: асимметрия стоимости (gold = «секунды → часы») ──
    ax, ay, aw, ah = 0.40, 1.94, 12.55, 2.02
    ocean_box(s, ax, ay, aw, ah)
    text_box(s, ax + 0.24, ay + 0.12, aw - 0.48, 0.30,
             "Асимметрия стоимости — почему это атака, а не просто спам",
             size=13.5, bold=True, color=MID)
    colw = (aw - 0.48 - 1.18) / 2.0
    cy = ay + 0.52
    cyh = ah - 0.68
    # LEFT — сгенерировать фейк: дёшево (нейтральный surface)
    lcx = ax + 0.24
    filled_rect(s, lcx, cy, colw, cyh, SURFACE, stroke=SOFT_GREY,
                stroke_pt=1.0, radius=True, radius_adj=0.06)
    text_box(s, lcx + 0.18, cy + 0.16, colw - 0.36, 0.32,
             "Сгенерировать фейк-отчёт", size=14, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, lcx + 0.18, cy + 0.52, colw - 0.36, 0.42,
             "секунды · почти ноль усилий", size=16, bold=True, color=MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lcx + 0.18, cy + 0.96, colw - 0.36, cyh - 1.08,
             "правдоподобный жаргон, дампы, ссылки — по форме не отличить "
             "от настоящего без ручной проверки",
             size=11.5, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
    # Center — gold asymmetry axis «секунды → часы» (× 1000+ = magnitude)
    # [Fix B1, issue #162] gold TEXT on white/transparent (≈2.0:1 FAIL) →
    # solid-gold chip behind each label + DEEP text (≈6.85:1 PASS).
    midx = lcx + colw + 0.14
    filled_rect(s, midx - 0.04, cy + cyh / 2 - 0.62, 0.98, 0.30, GOLD,
                radius=True, radius_adj=0.30)
    text_box(s, midx - 0.04, cy + cyh / 2 - 0.62, 0.98, 0.30,
             "× 1000+", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    right_arrow(s, midx, cy + cyh / 2 - 0.24, 0.90, 0.50, fill=GOLD)
    filled_rect(s, midx - 0.04, cy + cyh / 2 + 0.32, 0.98, 0.30, GOLD,
                radius=True, radius_adj=0.30)
    text_box(s, midx - 0.04, cy + cyh / 2 + 0.32, 0.98, 0.30,
             "разрыв", size=11.5, bold=True, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # RIGHT — опровергнуть: дорого (gold = ось)
    rcx = midx + 0.98
    filled_rect(s, rcx, cy, colw, cyh, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.06)
    text_box(s, rcx + 0.18, cy + 0.16, colw - 0.36, 0.32,
             "Опровергнуть фейк", size=14, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, rcx + 0.18, cy + 0.52, colw - 0.36, 0.42,
             "часы человека", size=16, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rcx + 0.18, cy + 0.96, colw - 0.36, cyh - 1.08,
             "понять сценарий · воспроизвести (или показать "
             "невоспроизводимость) · проверить смежный код",
             size=11.5, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
    # Conclusion band — DDoS на внимание / supply-chain (teal = системно)
    teal_callout(s, 0.55, 4.12, 12.25, 0.98,
                 "Одна сторона дешевеет в тысячи раз, другая нет — процесс "
                 "ломается. Атакуется не сервер, а невосполнимый ресурс — "
                 "время разбора сопровождающих (DDoS на внимание). Закроют "
                 "канал — реальные уязвимости утонут в шуме, риск уходит во "
                 "всю supply-chain.", size=13)
    # Lesson + non-AI alternative (gold callout, как s21/s22)
    gold_callout(s, 0.55, 5.24, 12.25, 1.04,
                 "Виновата архитектура процесса, не «AI» (запрет "
                 "неисполним). Чинить процесс, не-AI: приватное раскрытие "
                 "(GitHub Security Advisories) · убрать junk-стимул · "
                 "барьер воспроизводимого PoC на входе — машинный критерий "
                 "вместо человеческой проверки правдоподобия.", size=14)
    footer(s, "Критерий «когда AI здесь опасен»: открытый процесс, "
              "принимающий правдоподобный текст от кого угодно, И дорогая "
              "валидация на людях. Конкретные доли уточняются, асимметрия "
              "стабильна.")
    speaker_notes(s, load_notes("s22a"))


def build_s23(p):
    """case_study — 2 canons + CamoLeak + 4 Лекция-3 rules."""
    s = blank(p)
    slide_title(s, "Недоверенный ввод + привилегии = утечка.", size=26)
    # 2 canons
    c1x, cy, cw1, ch1 = 0.55, 1.40, 6.0, 2.05
    ocean_box(s, c1x, cy, cw1, ch1)
    text_box(s, c1x + 0.22, cy + 0.14, cw1 - 0.44, 0.34,
             "Канон 1 — корп.код/секреты в публичный чат = утечка",
             size=13.5, bold=True, color=MID, line_spacing=1.05)
    text_box(s, c1x + 0.22, cy + 0.56, cw1 - 0.44, 1.40,
             "данные за периметром живут по правилам, на которые вы не "
             "влияете (срок хранения, судебные приказы) — прямой перенос "
             "урока Лекции 3",
             size=12.5, color=DEEP, line_spacing=1.20)
    ocean_box(s, 6.80, cy, cw1, ch1, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, 7.02, cy + 0.14, cw1 - 0.44, 0.34,
             "Канон 2 — prompt-injection (CamoLeak)",
             size=13.5, bold=True, color=DEEP)
    text_box(s, 7.02, cy + 0.50, cw1 - 0.44, 1.50,
             "скрытые в невидимом тексте PR инструкции заставляли Copilot "
             "Chat искать AWS-ключи и отправлять наружу. Это «сбитый-с-толку "
             "посредник» (confused-deputy): агент исполняет чужую инструкцию "
             "своими правами — чужой PR стал командой.",
             size=11.5, color=DEEP, line_spacing=1.16)
    # 4 rules
    ocean_box(s, 0.55, 3.60, 12.25, 1.62)
    text_box(s, 0.78, 3.72, 11.8, 0.32,
             "Защита — архитектурная, ровно 4 правила Лекции 3:",
             size=14, bold=True, color=MID)
    rules = [
        ("least-privilege", "не давать широкий доступ к секретам/репо"),
        ("изоляция", "чужой PR/issue ≠ привилегированный контекст"),
        ("human-in-the-loop", "на write/деструктив"),
        ("egress-контроль", "ограничить, куда агент шлёт данные"),
    ]
    rxx = 0.78
    rcw = 5.95
    for i, (a, b) in enumerate(rules):
        col = i % 2
        row = i // 2
        xx = rxx + col * (rcw + 0.10)
        yy = 4.08 + row * 0.56
        text_box(s, xx, yy, 2.05, 0.46, a, size=12.5, bold=True, color=TEAL,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, xx + 2.08, yy, rcw - 2.20, 0.46, b, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    gold_callout(s, 0.55, 5.36, 12.25, 0.88,
                 "Структурное свойство, а не баг продукта. Лекция 3 это уже "
                 "доказала; CamoLeak — её dev-инстанс. Защита та же, не "
                 "новая.", size=14.5)
    footer(s, "CamoLeak — задокументированный prompt-injection в Copilot "
              "Chat (раскрыт 2025, исправлен вендором).")
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    """summary / schema_matrix — 4 return-points → risk → control."""
    s = blank(p)
    slide_title(s, "Каждый риск — и его конкретный контроль.", size=26)
    ocean_box(s, 0.40, 1.30, 12.55, 3.85)
    tx, ty = 0.55, 1.44
    col_w = [0.85, 5.10, 6.25]
    headers = ["№", "Риск", "Контроль (часто не-AI)"]
    hh, rh = 0.52, 0.78
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID)
        text_box(s, cx + 0.14, ty, col_w[j] - 0.24, hh, hd,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                 align=(PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT))
        cx += col_w[j]
    rows = [
        ("1", "«почти правильный» код",
         "человеческое ревью + тест до доверия фрагменту"),
        ("2", "merge без чтения, рост техдолга",
         "ревью PR + CI-gate (клоны / churn / mutation)"),
        ("3", "деструктив без гейта; перекос ощущения",
         "hard human-gate на необратимое; измерять, не ощущать"),
        ("4", "уязвимый код + ложная уверенность; утечка",
         "SAST/secret-scan; least-privilege + изоляция + egress"),
    ]
    yy = ty + hh
    for ri, (c0, c1, c2) in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else SURFACE
        cx = tx
        for j, cc in enumerate((c0, c1, c2)):
            filled_rect(s, cx, yy, col_w[j], rh, bg, stroke=SOFT_GREY,
                        stroke_pt=0.75)
            text_box(s, cx + 0.14, yy, col_w[j] - 0.24, rh, cc,
                     size=(15 if j == 0 else 12.5), bold=(j == 0),
                     color=(MID if j == 0 else DEEP),
                     align=(PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT),
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 5.32, 12.25, 0.94,
                 "Каждый риск имеет конкретный, известный, часто не-AI "
                 "контроль. Задача инженера — не бояться AI и не доверять "
                 "ему, а знать, какой контроль ставится и почему он не "
                 "делегируется.", size=14)
    speaker_notes(s, load_notes("s24"))


def build_s24a(p):
    """section_divider — Раздел 5: Методологии, конфигурации, люди.
    [Решение #101 owner GATE B] suffix-ID, cascade-safe; here_idx=5 →
    roadmap gold-маркер Раздел 5. Шаблон единый с s10/s14/s18."""
    build_section_divider(
        p, 5, "Методологии,\nконфигурации, люди",
        "Разделы 1–4 были про AI-код и его риски. Дальше — про процесс и "
        "людей: какие методологии ложатся на AI, чем solo+AI отличается от "
        "команды. AI меняет цену рутины — но методологии и роли не уходят, "
        "они уточняются.", "s24a")


def build_s25(p):
    """matrix / schema_matrix — methodologies × AI compatibility.
    WATCH-ITEM: peak density. axis-in, ≥14pt, ≥75% fill, solid color."""
    s = blank(p)
    slide_title(s, "Не все методологии одинаково совместимы с AI.",
                size=25, y=0.32, h=0.58)
    text_box(s, 0.55, 0.88, 12.25, 0.30,
             "Совместимость падает сверху вниз: TDD ложится на AI лучше "
             "всего, vibe-coding — антипаттерн.",
             size=12.5, italic=True, color=MID)
    ocean_box(s, 0.40, 1.24, 12.55, 3.55)
    tx, ty = 0.55, 1.38
    col_w = [3.30, 8.85]
    headers = ["Методология", "Почему ложится / что меняется"]
    hh = 0.54
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID)
        text_box(s, cx + 0.16, ty, col_w[j] - 0.30, hh, hd,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    rows = [
        ("TDD — №1",
         "тест пишется до кода = точная исполняемая спецификация + "
         "детерминированная обратная связь для агента (не самооценка)",
         "g"),
        ("spec-driven (спека до кода)",
         "контракт (требования / дизайн / задачи) сужает пространство "
         "правдоподобно-неверного", "n"),
        ("trunk-based + CI-гейты",
         "компенсирует системный эффект: AI ускоряет поток изменений, "
         "но снижает стабильность доставки (данные DORA)", "n"),
        ("vibe-coding — антипаттерн",
         "без структуры/ограничений/теста/гейта — снимает все 3 опоры "
         "инженерного паттерна; корень провалов лекции", "x"),
    ]
    rh = 0.74
    yy = ty + hh
    for nm, desc, sem in rows:
        cx = tx
        # gold SOLID = сильная сторона (TDD №1); teal SOLID = антипаттерн,
        # отвергается по построению (vibe-coding); surface = нейтральная.
        if sem == "g":
            c0bg, c0fg = GOLD, DEEP
            c1bg, c1fg = GOLD_TINT, DEEP
        elif sem == "x":
            c0bg, c0fg = TEAL, WHITE
            c1bg, c1fg = TEAL_TINT, DEEP
        else:
            c0bg, c0fg = SURFACE, DEEP
            c1bg, c1fg = WHITE, DEEP
        bold0 = sem in ("g", "x")
        filled_rect(s, cx, yy, col_w[0], rh, c0bg, stroke=WHITE,
                    stroke_pt=1.5)
        text_box(s, cx + 0.12, yy, col_w[0] - 0.22, rh, nm,
                 size=14, bold=bold0, color=c0fg, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        cx += col_w[0]
        filled_rect(s, cx, yy, col_w[1], rh, c1bg,
                    stroke=(GOLD if sem == "g" else
                            TEAL if sem == "x" else SOFT_GREY),
                    stroke_pt=1.2 if sem in ("g", "x") else 0.75)
        text_box(s, cx + 0.16, yy + 0.04, col_w[1] - 0.30, rh - 0.08, desc,
                 size=13, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.10)
        yy += rh
    gold_callout(s, 0.55, 4.92, 12.25, 1.34,
                 "vibe-coding (строгое определение): генерировать и "
                 "принимать код «по ощущению» — без явной структуры задачи, "
                 "без сформулированных ограничений, без машинно-проверяемого "
                 "теста и без quality-gate, доверяя правдоподобию вывода. "
                 "Это не методология, а её отсутствие.", size=14)
    footer(s, "DORA 2025: «AI — амплификатор, и TDD усиливается». Индустрия "
              "отвергает не AI, а паттерн vibe-coding-без-гейтов.")
    speaker_notes(s, load_notes("s25"))


def build_s26(p):
    """case_study — Brooks accidental/essential + DORA amplifier.
    Один ясный тезис для новичка: AI убирает рутину, не «решить что
    строить»; практики не уходят — уточняются (DORA: AI усиливает то,
    что уже есть). EN-цитаты → speaker notes."""
    s = blank(p)
    slide_title(s, "AI убирает рутину — но не решение, что строить.",
                size=24, y=0.34, h=0.58)
    text_box(s, 0.55, 0.96, 12.25, 0.34,
             "Классика инженерии (Брукс, 1986): сложность ПО бывает двух "
             "родов — и AI помогает только с одним.",
             size=13, italic=True, color=MID)
    # Two complexity kinds — parallel cards
    cy, chh = 1.46, 2.30
    cw2 = 6.0
    # accidental — AI силён
    ocean_box(s, 0.55, cy, cw2, chh, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, 0.79, cy + 0.16, cw2 - 0.48, 0.66,
             "Привнесённая сложность (accidental)",
             size=15, bold=True, color=TEAL, line_spacing=1.05)
    text_box(s, 0.79, cy + 0.78, cw2 - 0.48, 0.74,
             "трудность не от задачи, а от инструментов и рутины: "
             "boilerplate (шаблонный код), ручная отладка, документация",
             size=12.5, color=DEEP, line_spacing=1.16)
    text_box(s, 0.79, cy + 1.56, cw2 - 0.48, 0.62,
             "→ здесь AI силён: эту рутину он реально удешевляет",
             size=13, bold=True, color=DEEP, line_spacing=1.12)
    # essential — AI не помогает (gold = главный тезис)
    ocean_box(s, 6.80, cy, cw2, chh, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, 7.04, cy + 0.16, cw2 - 0.48, 0.66,
             "Существенная сложность (essential)",
             size=15, bold=True, color=DEEP, line_spacing=1.05)
    text_box(s, 7.04, cy + 0.78, cw2 - 0.48, 0.74,
             "трудность самой задачи: точно решить, ЧТО именно "
             "строить, — выбор, постановка, ответственность",
             size=12.5, color=DEEP, line_spacing=1.16)
    text_box(s, 7.04, cy + 1.56, cw2 - 0.48, 0.62,
             "→ здесь AI не помогает: это и есть «что не делегируется»",
             size=13, bold=True, color=DEEP, line_spacing=1.12)
    # DORA strip — practices don't go away [Fix C, #162: + ROI-отчёт май
    # 2026, $-квантификация + J-кривая, не только качественный вывод 2025]
    ocean_box(s, 0.55, 3.94, 12.25, 1.42)
    text_box(s, 0.79, 4.04, 11.7, 0.30,
             "DORA 2025 (~5000 команд) + ROI-отчёт (май 2026) — почему "
             "практики не уходят:",
             size=13.5, bold=True, color=MID)
    text_box(s, 0.79, 4.36, 11.7, 0.42,
             "AI внедрили ~90% команд, но связь AI со стабильностью "
             "доставки — отрицательная, второй год подряд: «AI не чинит "
             "команду — он усиливает то, что уже есть».",
             size=12.5, color=DEEP, line_spacing=1.16)
    text_box(s, 0.79, 4.82, 11.7, 0.42,
             "ROI-отчёт (2026): та же условность количественно — рост доли "
             "неудачных изменений 5%→6% ≈ −$344 000; J-кривая — временный "
             "провал стабильности без готовых gate заранее.",
             size=12.5, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 5.48, 12.25, 0.94,
                 "AI — усилитель, не исправитель. Сильную инженерную "
                 "культуру (тесты, ревью, гейты) AI делает сильнее; "
                 "слабую — ломает быстрее, и это теперь измеримо в "
                 "долларах, а не только в абстрактной «стабильности».",
                 size=13.5)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """comparison / schema_matrix — solo+AI vs team+AI tradeoff."""
    s = blank(p)
    slide_title(s, "Конфигурация — это размен под задачу, не «что круче».",
                size=25, y=0.34, h=0.58)
    # 2 equal columns
    bx, by, bw, bh = 0.55, 1.36, 12.25, 2.30
    ocean_box(s, bx, by, bw, bh)
    cw2 = (bw - 0.60) / 2
    # solo
    sx = bx + 0.20
    filled_rect(s, sx, by + 0.18, cw2, bh - 0.36, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.05)
    text_box(s, sx + 0.22, by + 0.32, cw2 - 0.44, 0.38, "solo + AI",
             size=17, bold=True, color=DEEP)
    text_box(s, sx + 0.22, by + 0.78, cw2 - 0.44, 0.56,
             "дёшево / быстро, нет затрат на согласование между людьми",
             size=12.5, color=DEEP, line_spacing=1.14)
    text_box(s, sx + 0.22, by + 1.36, cw2 - 0.44, 0.78,
             "истощённое узкое место: один человек 24/7, нет второго "
             "ревьюера → единственная точка отказа",
             size=12.5, bold=True, color=DEEP, line_spacing=1.16)
    # team
    tx2 = bx + 0.40 + cw2
    filled_rect(s, tx2, by + 0.18, cw2, bh - 0.36, SURFACE, stroke=SOFT_GREY,
                stroke_pt=1.0, radius=True, radius_adj=0.05)
    text_box(s, tx2 + 0.22, by + 0.32, cw2 - 0.44, 0.38, "команда + AI",
             size=17, bold=True, color=MID)
    text_box(s, tx2 + 0.22, by + 0.78, cw2 - 0.44, 0.56,
             "взаимное ревью, распределённая ответственность, владение "
             "кодом", size=12.5, color=DEEP, line_spacing=1.14)
    text_box(s, tx2 + 0.22, by + 1.36, cw2 - 0.44, 0.78,
             "AI усиливает сильную команду (а слабую — ломает быстрее, "
             "данные DORA)",
             size=12.5, color=DEEP, line_spacing=1.16)
    # criterion
    ocean_box(s, 0.55, 3.78, 6.0, 1.55)
    # [Fix B1, issue #162] gold TEXT on SURFACE (≈1.85:1 FAIL) → solid-gold
    # chip + DEEP text (≈6.85:1 PASS).
    filled_rect(s, 0.78, 3.88, 1.62, 0.34, GOLD, radius=True, radius_adj=0.25)
    text_box(s, 0.78, 3.88, 1.62, 0.34, "solo + AI", size=13, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.78, 4.22, 5.55, 1.05,
             "ранняя стадия / прототип / узкая чёткая задача / "
             "обратимые последствия",
             size=12.5, color=DEEP, line_spacing=1.18)
    ocean_box(s, 6.80, 3.78, 6.0, 1.55)
    text_box(s, 7.03, 3.90, 5.55, 0.30, "команда + AI", size=13, bold=True,
             color=MID)
    text_box(s, 7.03, 4.22, 5.55, 1.05,
             "прод / регулируемое / долгоживущий код / нужен аудит и "
             "распределённая ответственность",
             size=12.5, color=DEEP, line_spacing=1.18)
    teal_callout(s, 0.55, 5.46, 12.25, 0.80,
                 "Ландшафт (направление, не точные доли): Copilot стагнирует "
                 "(всё ещё #1 по охвату); Claude Code / Cursor растут; "
                 "уходит не инструмент, а практика vibe-coding-без-гейтов.",
                 size=13)
    footer(s, "Человеческое суждение (что строить, цена, рынок) — "
              "невосполнимое ядро в ОБЕИХ конфигурациях. Доли охвата "
              "инструментов меняются почти еженедельно — важно направление, "
              "не точное число.")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """assertion_visual — docs-as-code: зачем для AI важнее, чем для
    человека (left, подтверждено) vs где слабое место (right, vendor-
    claim). Новичку: одна фраза-определение ДО использования."""
    s = blank(p)
    slide_title(s, "Для AI документация важнее, чем для человека — но не вся.",
                size=24, y=0.34, h=0.58)
    text_box(s, 0.55, 0.96, 12.25, 0.40,
             "docs-as-code — документация лежит в репозитории рядом с "
             "кодом, под версионным контролем и в том же ревью, что и код.",
             size=13, italic=True, color=MID, line_spacing=1.16)
    # Left — ЗАЧЕМ для AI важнее (подтверждено)
    z1x, zy, zw, zh = 0.55, 1.62, 6.0, 3.50
    ocean_box(s, z1x, zy, zw, zh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, z1x + 0.24, zy + 0.16, zw - 0.48, 0.62,
             "Почему для AI важнее (подтверждено)",
             size=14.5, bold=True, color=TEAL, line_spacing=1.05)
    conf = [
        "человек добирает контекст сам — догадается, спросит коллегу; "
        "агент видит только то, что ему дали текстом",
        "AGENTS.md / CLAUDE.md — файл с контекстом и правилами для "
        "агента — де-факто стандарт (с августа 2025)",
        "рост с ~20k до 40k+ репозиториев, нативная поддержка в "
        "инструментах — это документация для машины, не для людей",
    ]
    cyy = zy + 0.86
    for t in conf:
        circle(s, z1x + 0.26, cyy + 0.06, 0.12, TEAL)
        text_box(s, z1x + 0.50, cyy, zw - 0.74, 1.30, t, size=12,
                 color=DEEP, line_spacing=1.14)
        cyy += 0.90
    # Right — где слабое место (vendor-claim)
    z2x = 6.80
    ocean_box(s, z2x, zy, zw, zh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, z2x + 0.24, zy + 0.16, zw - 0.48, 0.62,
             "Где слабое место (пока не доказано)",
             size=14.5, bold=True, color=DEEP, line_spacing=1.05)
    weak = [
        "тезис «спецификация замещает код как единственный источник "
        "истины» — пока заявление вендоров, не независимое исследование",
        "сами же эти практики признают: источник истины остаётся код",
        "цифры ускорения в 3–10× — отчёты ранних адоптеров, не "
        "проверенные данные",
    ]
    wyy = zy + 0.86
    for t in weak:
        circle(s, z2x + 0.26, wyy + 0.06, 0.12, GOLD)
        text_box(s, z2x + 0.50, wyy, zw - 0.74, 1.30, t, size=12,
                 color=DEEP, line_spacing=1.14)
        wyy += 0.90
    gold_callout(s, 0.55, 5.30, 12.25, 0.96,
                 "Контекст для агента в репозитории — использовать, это "
                 "работает. «Спека вместо кода как истина» — не закладывать "
                 "в архитектуру как факт, пока нет независимого "
                 "подтверждения.", size=14)
    footer(s, "Числа охвата AGENTS.md/CLAUDE.md быстро растут — важен факт "
              "стандарта, не точное число. «Спека = единственная истина» — "
              "заявление вендора, помечаем «слабо подтверждено».")
    speaker_notes(s, load_notes("s28"))


def build_s28a(p):
    """section_divider — Раздел 6: Фреймворк решения. [Решение #101 owner
    GATE B] suffix-ID, cascade-safe; here_idx=6 → roadmap gold-маркер
    Раздел 6. Шаблон единый с s10/s14/s18."""
    build_section_divider(
        p, 6, "Фреймворк решения",
        "Всё, что разобрали — «где ускоряет / где вредит / что не "
        "делегируется» — собираем в один аппарат: матрица уровень × задача, "
        "критерий «когда не и опасно», чек-лист перед тем, как дать задачу "
        "AI.", "s28a")


def build_s29(p):
    """matrix → LO7 payoff: доминанта-вывод + 5 канонических осей выбора
    ПО ИМЕНАМ компактной лентой (chapter §6.1), читаемые ≥14pt из зала.
    «Решающая ось, не сумма баллов». Полная сетка 5×4 → speaker notes.
    «Повторяемость» — НЕ ось матрицы (это pre-фильтр «не AI», нижняя
    плашка). Читаемость на 68-й мин — приоритет."""
    s = blank(p)
    slide_title(s, "Чем мерить требование задачи — пять осей выбора.",
                size=25, y=0.32, h=0.56)
    # DOMINANT conclusion plate (the synthesized takeaway — 5-sec anchor)
    dx, dy, dw, dh = 0.55, 1.00, 12.25, 1.30
    filled_rect(s, dx, dy, dw, dh, GOLD_TINT, stroke=GOLD, stroke_pt=3.0,
                radius=True, radius_adj=0.07)
    icon(s, "scale", dx + 0.32, dy + (dh - 0.58) / 2, 0.58, "gold")
    text_box(s, dx + 1.12, dy + 0.14, dw - 1.42, dh - 0.28,
             "Простое и обратимое AI ведёт сам; сложное, незнакомое и "
             "необратимое — руки + человеческий гейт. Решает не сумма "
             "баллов, а главная ось задачи.",
             size=17, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16)
    # 5 canonical axes — named ribbon, ≥14pt names (chapter §6.1)
    text_box(s, 0.55, 2.46, 12.25, 0.28,
             "Пять осей — для каждой задачи назвать решающую, проверить "
             "остальные на блокеры:",
             size=13, bold=True, color=MID)
    axes = [
        ("git-compare", "Незнакомость кода",
         "знакомый типовой → AI надёжен · приватный/легаси → надёжность падает"),
        ("refresh-cw", "Обратимость операции",
         "обратимо → выше автономия · необратимо → жёсткий человеческий гейт"),
        ("flame", "Критичность / прод",
         "некритичное → шире автономия · прод и пользователи → строгий гейт"),
        ("gavel", "Аудит / ответственность",
         "нет следа → можно · нужен владелец решения → человек на merge"),
        ("scale", "Цена ошибки",
         "задаёт приоритет осей — она вето: одна высокая опускает потолок"),
    ]
    n = len(axes)
    rx0, ry0 = 0.40, 2.80
    rw, rgap = 12.55, 0.10
    rh = 0.50
    for i, (ic, ttl, sub) in enumerate(axes):
        yy = ry0 + i * (rh + rgap)
        is_price = (ttl == "Цена ошибки")
        filled_rect(s, rx0, yy, rw, rh,
                    GOLD_TINT if is_price else SURFACE,
                    stroke=GOLD if is_price else SOFT_GREY,
                    stroke_pt=2.0 if is_price else 1.0,
                    radius=True, radius_adj=0.16)
        icon(s, ic, rx0 + 0.16, yy + (rh - 0.32) / 2, 0.32,
             "gold" if is_price else "mid")
        text_box(s, rx0 + 0.62, yy, 3.35, rh, ttl,
                 size=14, bold=True, color=DEEP if is_price else MID,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx0 + 4.05, yy, rw - 4.20, rh, sub,
                 size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    # bottom anchor — pre-filter «not AI at all» (where «повторяемость» lives)
    teal_callout(s, 0.55, 5.92, 12.25, 0.74,
                 "Сначала отсев: детерминированная, проверяемая, повторяемая "
                 "задача (парсинг, валидация по схеме, арифметика) → обычный "
                 "код, без AI вовсе — остальные оси не нужны.",
                 size=13, bold=False)
    footer(s, "Синтез разделов: METR · Brooks · 70/80%-проблема · разрыв "
              "SWE-bench Verified↔Pro · Replit/CamoLeak.")
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """assertion_visual — 4 classes where lower autonomy / not AI."""
    s = blank(p)
    slide_title(s, "Когда правильный ответ — понизить автономию или не AI.",
                size=25, y=0.34, h=0.58)
    cards = [
        ("circle-slash", "1. Детерминированная верифицируемая → не AI вовсе",
         "парсинг, валидация по схеме, арифметика, маршрутизация — обычный "
         "код точен и аудируем", "mid", False),
        ("bomb", "2. High-stakes без ревью → недопустимо",
         "необратимое/критичное/прод без человеческого гейта — профиль "
         "Replit/Kiro/PocketOS", "mid", False),
        ("graduation-cap", "3. Обучение junior делегированием → вредит навыку",
         "Anthropic RCT: делегировавшие — −17% на квизе; спрашивавшие "
         "«как работает» — без деградации", "gold", True),
        ("circle-x", "4. Автономия без hard-гейта на необратимое → запрещено",
         "уровень D без least-privilege / rollback / egress — не "
         "настраиваемый параметр", "mid", False),
    ]
    cw2 = 6.0
    ch2 = 1.78
    for i, (ic, ttl, body, var, hi) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = 0.55 + col * (cw2 + 0.25)
        cy = 1.34 + row * (ch2 + 0.20)
        if hi:
            ocean_box(s, cx, cy, cw2, ch2, fill=GOLD_TINT, stroke=GOLD,
                      stroke_pt=2.0)
        else:
            ocean_box(s, cx, cy, cw2, ch2)
        icon(s, ic, cx + 0.24, cy + 0.22, 0.46, var)
        text_box(s, cx + 0.84, cy + 0.20, cw2 - 1.05, 0.74, ttl,
                 size=13.5, bold=True, color=(DEEP if hi else MID),
                 line_spacing=1.10, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + 0.26, cy + 1.00, cw2 - 0.52, ch2 - 1.14, body,
                 size=12, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 5.30, 12.25, 0.96,
                 "Глава не «бойтесь AI» и не «AI всё решит». Между "
                 "AI-карго-культом и AI-отрицанием — инженер, который для "
                 "каждой задачи называет уровень, конфигурацию, точку "
                 "человеческого контроля и условие смены.", size=14)
    footer(s, "Anthropic «How AI Impacts Skill Formation» (2026, n=52): "
              "−17% / >60% делегировали / ≥65% спрашивали концепции.")
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """summary — 8-step checklist + worked example + mini-apply."""
    s = blank(p)
    slide_title(s, "Чек-лист «прежде чем дать задачу AI».", size=26)
    # Left — checklist
    lx, ly, lw, lh = 0.55, 1.34, 7.05, 5.50
    ocean_box(s, lx, ly, lw, lh)
    items = [
        "Можно ли решить без AI (детерм., верифиц.)? → не добавляй AI",
        "Обратимо ли последствие? Необратимое → hard human-gate",
        "Есть ли тест-оракул? Нет → ревью обязательно",
        "Кто ревьюит и кто мержит? Merge — всегда человек",
        "Затронуты секреты / недоверенный контент? → least-priv + изоляция",
        "Насколько код знаком AI? Незнакомый → доверие по Pro (~69–80%)",
        "Цель — артефакт или навык? Навык → не делегировать генерацию",
        "Конфигурация: solo+AI или команда+AI — по обратимости/аудиту?",
    ]
    iy = ly + 0.26
    for i, it in enumerate(items):
        circle(s, lx + 0.26, iy + 0.04, 0.34, MID)
        text_box(s, lx + 0.26, iy + 0.04, 0.34, 0.34, str(i + 1),
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, lx + 0.72, iy, lw - 0.98, 0.60, it,
                 size=12.5, color=DEEP, line_spacing=1.10,
                 anchor=MSO_ANCHOR.MIDDLE)
        iy += 0.635
    # Right — apply
    rx, rw = 7.85, 4.95
    ocean_box(s, rx, 1.34, rw, 2.55, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, rx + 0.24, 1.46, rw - 0.48, 0.32,
             "Worked example (задача A)", size=13, bold=True, color=TEAL)
    text_box(s, rx + 0.24, 1.80, rw - 0.48, 2.00,
             "приватный платёжный модуль, тесты неполные → C с "
             "обязательным senior-ревью + hard-гейт, команда+AI. Решающая "
             "ось — необратимость+критичность. Условие смены: обратимый "
             "внутренний инструмент с полным покрытием → допустим D.",
             size=11.5, color=DEEP, line_spacing=1.16)
    ocean_box(s, rx, 4.04, rw, 1.55)
    text_box(s, rx + 0.24, 4.16, rw - 0.48, 0.32,
             "Mini-apply (задача B)", size=13, bold=True, color=MID)
    text_box(s, rx + 0.24, 4.50, rw - 0.48, 1.00,
             "миграция формата конфигов в ~300 репо по фикс-правилу, "
             "мержит человек — пройдите чек-лист сами за 2 минуты "
             "(think-pair-share)",
             size=11.5, color=DEEP, line_spacing=1.16)
    gold_callout(s, 7.85, 5.74, 4.95, 1.10,
                 "Сначала попытка — потом сверка. Сформулируйте ответ ДО "
                 "разбора. Полная отработка — Семинар 4.", size=12)
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    """hero_closing / qa_minimal — bridge to industry lectures + Семинар 4
    + Q&A. HERO (Fix B2, issue #162): real screenshot of the Replit
    agent's own confession chat log (s16 — most memorable failure artefact
    of the lecture), Tier 5 acquisition (Wayback Machine), ≥40% slide area
    — closes the emotional arc «AI can be confidently, articulately wrong
    about its own actions»."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    # ---- LEFT column: bridge + homework + Q&A (compact) ----
    bx, by, bw = 0.55, 0.55, 6.35
    ocean_box(s, bx, by, bw, 1.42)
    icon(s, "route", bx + 0.22, by + 0.26, 0.44, "mid")
    text_box(s, bx + 0.78, by + 0.14, bw - 0.98, 0.38,
             "Лекция 4 — первая отраслевая",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, bx + 0.78, by + 0.54, bw - 0.98, 0.78,
             "Лестница A→D + матрица выбора + чек-лист — линза для всех "
             "следующих отраслевых лекций: тот же вопрос в новой обёртке.",
             size=12, color=DEEP, line_spacing=1.14)
    # homework block
    ocean_box(s, bx, by + 1.58, bw, 1.62, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, bx + 0.26, by + 1.70, bw - 0.52, 0.56,
             "Задание — Семинар 4: «AI в цикле разработки ПО: "
             "автодополнение, чат-ассистент, агент»",
             size=13, bold=True, color=DEEP, line_spacing=1.08)
    text_box(s, bx + 0.26, by + 2.28, bw - 0.52, 0.84,
             "Для каждого кейса: уровень + ≥2 причины по осям + ≥1 условие "
             "смены + явное «где человек обязателен». Mini-apply задача B "
             "— разминка.",
             size=11.5, color=DEEP, line_spacing=1.14)
    # Q&A (compact, left column)
    text_box(s, bx, by + 3.42, bw, 1.35, "Вопросы",
             size=56, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, bx, by + 4.72, bw, 0.44, "Спасибо за внимание",
             size=18, bold=False, color=MID, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, bx, by + 5.20, bw, 0.60,
             "Семинар 4 — на следующей неделе. Дополнительные вопросы — "
             "на e-mail.", size=11.5, italic=True, color=LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    # ---- RIGHT column: HERO — Replit agent confession screenshot ----
    rx, rw = 7.15, 5.65
    ocean_box(s, rx, by, rw, 5.60, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, str(ROOT / "assets/screenshots/s32-replit-real-source.jpg"),
              x=rx + 0.14, y=by + 0.14, w=rw - 0.28, h=4.90)
    text_box(s, rx + 0.14, by + 5.08, rw - 0.28, 0.24,
             "X (Jason Lemkin) · 18 июля 2025 · архив Wayback Machine",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    text_box(s, rx + 0.14, by + 5.30, rw - 0.28, 0.28,
             "Агент подтверждает: удалил прод-БД в code-freeze",
             size=11.5, bold=True, italic=True, color=DEEP,
             align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s32"))


# ============================================================
# Main
# ============================================================
def main():
    spec = load_deck()
    if spec is not None:
        n = len(spec["slides"])
        print(f"deck spec OK — {n} slides (deck.yaml + deck-part2.yaml), "
              f"totals {spec['totals'].get('slides')}")
    p = setup_pres()
    # [Решение #101, 2026-05-17 — owner GATE B] 32 base (s01–s32 нумерация
    # неизменна) + 3 suffix-ID раздела-дивайдера: s04a после s04 (Р1),
    # s24a после s24 (Р5), s28a после s28 (Р6). cascade-safe.
    # [Решение #103, 2026-05-17 — owner GATE C] +1 контент-слайд s22a
    # (curl-slop #5) после s22, перед s23 (security-раздел). cascade-safe.
    builders = [build_s01, build_s02, build_s03, build_s04, build_s04a,
                build_s05, build_s06, build_s07, build_s08, build_s09,
                build_s10, build_s11, build_s12, build_s13,
                build_s13a, build_s13b, build_s13c, build_s13d, build_s13e,
                build_s13f, build_s13g, build_s13h,
                build_s14,
                build_s15, build_s16, build_s17, build_s18, build_s19,
                build_s20, build_s21, build_s22, build_s22a, build_s23,
                build_s24, build_s24a, build_s25, build_s26, build_s27,
                build_s28, build_s28a, build_s29, build_s30, build_s31,
                build_s32]
    assert len(builders) == 44, f"expected 44 builders, got {len(builders)}"
    for b in builders:
        b(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"saved {OUT} — {len(p.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
