"""
EN track — full 55-slide build of Lecture 3 "AI-system architectures:
agents, RAG, API" (v6.3-en, issue #172/#185).

Ported 1:1 from build_v3.py (RU source of truth, 55 slides). Same structure,
same builders, same layout — only the visible text + speaker notes are in
English. English memes (*-en.png) via make_memes_en.py. Anchors via
refs_lec03_en.py.

Source-of-truth: RU build_v3.py + deck.en.yaml + slides-en/*.md
(readable speaker notes 150-300 words).

Palette LOCKED: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) >=1x/slide.
Visual motif: "Ocean rounded box" (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).

Canvas: 13.333" x 7.5" (16:9).

Build via: python3 build_lec03_en.py — generates lec-03-en.pptx (55 slides).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree
from PIL import Image

import sys as _sys
_sys.path.insert(0, str(Path(__file__).resolve().parent))
import refs_lec03_en as R  # noqa: E402  (issue #172 EN reference/page-number system)

# issue #171: footer text capture so ref-slides can fold the caveat into the
# clickable [N] source list (single bottom band, no overlap).
_FOOTER_TEXT = {}

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x5B, 0x66, 0x78)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFD, 0xF3, 0xDC)
TEAL_TINT = RGBColor(0xE4, 0xF1, 0xF2)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
# issue #171: run from the worktree — read slides/assets/notes + write pptx
# from THIS repo checkout (slides identical to main). Falls back to main-repo
# ROOT only if the worktree copy is missing.
_WT = Path(__file__).resolve().parents[1]      # …/lec-03 in the current checkout
_MAIN = Path("/home/harness/harness-projects/256/lessons-3bb49d40/library/lectures/lec-03")
ROOT = _WT if (_WT / "slides").exists() else _MAIN
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
CHARTS = ASSETS / "charts-en"
SLIDES_DIR = ROOT / "slides-en"
OUT = ROOT / "rendered/lec-03-en.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before") is not None:
                p.space_before = Pt(cfg["space_before"])
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.035, min(0.22, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def up_arrow(slide, x, y, w, h, fill=LIGHT):
    shp = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, fill, *, stroke=None, stroke_pt=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def connector(slide, x1, y1, x2, y2, color=LIGHT, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = etree.SubElement(ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
        pd.set("val", dash)
    return cn


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    path = Path(path)
    if not path.exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path); iw, ih = img.size; img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
            return
        ir = iw / ih; br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x), Inches(y + (h - ah) / 2),
                                     width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2), Inches(y),
                                     height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


SCREENSHOTS = ROOT / "assets/screenshots"
_CROP_CACHE = ROOT / "rendered/assets/_crop_cache"


def hero_image(slide, src, x, y, w, h):
    """Cover-crop an image to EXACTLY fill box (x,y,w,h) — no distortion, no
    letterbox. python-pptx can't crop, so we pre-crop via PIL to the target
    aspect ratio and cache the result. Used for ≥40% hero fills (s01, s30)."""
    src = Path(src)
    if not src.exists():
        return
    _CROP_CACHE.mkdir(parents=True, exist_ok=True)
    target_ratio = w / h
    try:
        img = Image.open(src).convert("RGB")
    except Exception:
        return
    iw, ih = img.size
    ir = iw / ih
    if ir > target_ratio:      # image wider — crop sides
        new_w = int(ih * target_ratio)
        left = (iw - new_w) // 2
        img = img.crop((left, 0, left + new_w, ih))
    else:                       # image taller — crop top/bottom
        new_h = int(iw / target_ratio)
        top = (ih - new_h) // 2
        img = img.crop((0, top, iw, top + new_h))
    out = _CROP_CACHE / f"{src.stem}_{w:.2f}x{h:.2f}.png"
    img.save(out)
    slide.shapes.add_picture(str(out), Inches(x), Inches(y),
                             width=Inches(w), height=Inches(h))


def slide_title(slide, text, *, y=0.42, h=0.95, w=12.25, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.12, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 color=DEEP, align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.22, y=y + 0.06, w=w - 0.44, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.22)


def footer(slide, text):
    # issue #171: record the footer so ref-slides can relocate/fold it; still
    # render normally (post-processing removes it only on ref-slides).
    _FOOTER_TEXT[id(slide)] = text
    tb = text_box(slide, x=0.55, y=7.02, w=12.25, h=0.36, text=text,
                  size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
                  line_spacing=1.0)
    _FOOTER_TEXT.setdefault("_shapes", {})[id(slide)] = tb
    return tb


def icon(slide, name, x, y, size, variant="mid"):
    add_image(slide, ICONS / f"{name}-{variant}.png", x, y, size, size)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Deck loader — reads the multi-part RU deck.yaml for validation only.
# Merges the `slides` key across parts, validates totals (report-only).
# ============================================================
def load_deck():
    """Load 2-part deck spec (deck.yaml + deck-part2.yaml), merge `slides`.

    Returns dict with merged slide list + validation. Raises if the
    ordered slide-id list does not match the canonical v3 presentation
    order (cascade-safe guard — s01–s30 must NOT be renumbered).
    """
    try:
        import yaml
    except ImportError:
        return None  # yaml optional — builder list is authoritative anyway
    p1 = ROOT / "deck.yaml"
    p2 = ROOT / "deck-part2.yaml"
    parts = [ROOT / "deck.yaml", ROOT / "deck-part2.yaml", ROOT / "deck-part3.yaml"]
    slides = []
    d1 = None
    totals = {}
    for pp in parts:
        if not pp.exists():
            continue
        d = yaml.safe_load(pp.read_text(encoding="utf-8"))
        if d1 is None:
            d1 = d
        slides += list(d.get("slides", []))
        if d.get("totals"):
            totals = d["totals"]
    spec = {
        "deck": d1.get("deck", {}) if d1 else {},
        "palette": d1.get("palette", {}) if d1 else {},
        "slides": slides,
        "totals": totals,
    }
    ids = [s["id"] for s in slides]
    # builder list is authoritative; deck.yaml is documentation — report only.
    print(f"deck.yaml slide ids ({len(ids)}): {ids}")
    return spec


# ============================================================
# Section divider — unified template (6-card roadmap, gold current)
# Sections of Lecture 3 (deck.yaml): 0..5.
# ============================================================
NAV = [
    ("0", "Opening"),
    ("1", "Prompt"),
    ("2", "RAG"),
    ("3", "Fine-tuning"),
    ("4", "Agents"),
    ("5", "Framework"),
]


def roadmap_bar(slide, here_idx, *, y=6.45):
    """6-card progress bar; current section gold-bordered."""
    n = len(NAV)
    gap = 0.14
    bx = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    ch = 0.62
    for i, (num, label) in enumerate(NAV):
        x = bx + i * (cw + gap)
        cur = (i == here_idx)
        if cur:
            filled_rect(slide, x, y, cw, ch, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.12)
        else:
            filled_rect(slide, x, y, cw, ch, SURFACE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(slide, x=x + 0.06, y=y + 0.07, w=cw - 0.12, h=0.24,
                 text=f"Section {num}", size=10.5, bold=True,
                 color=(DEEP if cur else LIGHT), align=PP_ALIGN.CENTER)
        text_box(slide, x=x + 0.06, y=y + 0.31, w=cw - 0.12, h=0.26,
                 text=label, size=11, bold=cur,
                 color=(DEEP if cur else SLATE), align=PP_ALIGN.CENTER)


WEB = ASSETS / "web"


def build_section_divider(p, here_idx, big_num, subtitle, frame_phrase, sid,
                          *, image_src=None, image_caption=None, tag=None):
    """Distinct divider (NO ocean motif). Left = SECTION N + subtitle + frame
    phrase + optional tag chip; right = a REAL metaphor image (≥40% width,
    full height, cover-cropped) when image_src is given, else the giant
    cover-style decorative digit. Gold-current roadmap bar at bottom."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    have_img = image_src is not None and Path(image_src).exists()
    if have_img:
        # Real metaphor image — right ~42% width, full height (≥40% area).
        ix, iy, iw, ih = 8.10, 0.0, 5.233, 7.5
        hero_image(s, image_src, ix, iy, iw, ih)
        # faint decorative digit overlaid top-right of image (brand echo)
        text_box(s, x=ix + 1.6, y=0.20, w=3.6, h=2.6, text=str(here_idx),
                 size=170, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        if image_caption:
            filled_rect(s, ix, 7.10, iw, 0.40, DEEP)
            text_box(s, ix + 0.16, 7.12, iw - 0.30, 0.34, image_caption,
                     size=10, italic=True, color=WHITE,
                     anchor=MSO_ANCHOR.MIDDLE)
        left_w = 7.35
    else:
        text_box(s, x=8.35, y=0.55, w=4.6, h=5.6, text=str(here_idx),
                 size=380, bold=True, color=COVER_OUTLINE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        left_w = 7.5
    # Left — section label + subtitle + frame phrase
    text_box(s, x=0.75, y=1.35, w=left_w - 0.2, h=0.55,
             text=f"SECTION {here_idx}", size=20, bold=True, color=TEAL)
    filled_rect(s, 0.78, 1.96, 0.7, 0.05, fill=GOLD)
    text_box(s, x=0.75, y=2.30, w=left_w, h=1.75, text=subtitle,
             size=36, bold=True, color=DEEP, line_spacing=1.06)
    text_box(s, x=0.78, y=4.20, w=left_w - 0.15, h=1.45, text=frame_phrase,
             size=17, italic=True, color=LIGHT, line_spacing=1.20)
    if tag:
        chip(s, 0.78, 5.70, min(left_w - 0.3, 0.14 * len(tag) + 0.6), 0.44,
             tag, fill=TEAL, color=WHITE, size=12.5)
    roadmap_bar(s, here_idx, y=6.45)
    speaker_notes(s, load_notes(sid))
    return s


# ============================================================
# Slide builders — 30 slides
# ============================================================

def _pill_figure(s, cx, cy, scale, body_col, ok=True):
    """Flat vector figure holding up a giant 'pill'. ok=True → confident pose
    (gold pill, teal glow); ok=False → pill shattered / red-adjacent avoided,
    we keep Ocean palette: dull grey pill + slate figure."""
    # head
    circle(s, cx - 0.30 * scale, cy, 0.60 * scale, body_col)
    # torso (rounded rect)
    filled_rect(s, cx - 0.42 * scale, cy + 0.66 * scale, 0.84 * scale,
                1.15 * scale, body_col, radius=True, radius_adj=0.35)


def build_s01(p):
    """hero_cover / meme-hook — "magic pill" (issue #185). Real Drake
    internet meme (imgflip) with English captions: reject = complicate the
    prompt for accuracy, approve = pick the architecture for the task. The meme
    carries the thesis; the answer unfolds through the lecture. Air Canada — separately on s01b."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # top ribbon (minimum text — the framing in one line)
    text_box(s, 0.55, 0.40, 12.25, 0.44, "A MYTH ABOUT AI SYSTEMS",
             size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(s, 0.55, 0.86, 12.25, 0.92,
             "“Magic pill”: the model answers more accurately if you ask it to act as an expert and complicate the prompt.",
             size=25, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             line_spacing=1.06)
    # real Drake meme (hero ≥40%) in an Ocean frame, centered: on top we reject
    # "complicate the prompt for accuracy", below we accept "picking the architecture".
    my, mh = 2.06, 4.06
    mw = 4.94
    mx = (13.33 - mw) / 2
    ocean_box(s, mx - 0.16, my - 0.14, mw + 0.32, mh + 0.28)
    add_image(s, WEB / "s01-drake-en.png", mx, my, mw, mh)
    gold_callout(s, 0.55, 6.34, 12.25, 0.74,
                 "Where an AI system's reliability actually comes from — we work through it across the whole lecture.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """cover — distinct, NO ocean motif. Mega "03" + title + roadmap.
    v2: subtitle brought to lec-02 cover canon — content-promise line with
    teal accent bar + MID color (was designer-initiative "Course · 75 min"
    meta-line italic-light, removed)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=7.6, y=1.15, w=5.7, h=5.0, text="03",
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.75, y=1.35, w=6.6, h=0.5, text="LECTURE 3",
             size=18, bold=True, color=TEAL)
    filled_rect(s, 0.78, 1.92, 0.7, 0.05, fill=TEAL)
    text_box(s, x=0.75, y=2.35, w=7.7, h=2.7,
             text="AI system architectures:\nagents, RAG, API",
             size=46, bold=True, color=DEEP, line_spacing=1.08)
    # subtitle = content promise (lec-02 canon: teal accent bar + MID, not meta)
    filled_rect(s, 0.78, 5.28, 0.05, 0.56, fill=TEAL)
    text_box(s, x=1.02, y=5.26, w=7.4, h=0.62,
             text="Which architecture to pick for the task —\nand when the right answer is “not AI”",
             size=19, italic=False, color=MID, line_spacing=1.18)
    # v4 (#212): cover is now clean — roadmap moved to a separate
    # lecture-map slide s02a (L1/L2 pattern). Cover has no roadmap-bar.
    # #185/#313: course attribution "Year 3 IU6 · Module 1…" removed from cover.
    # #185/#313: a thematic internet meme on the cover (bottom-left corner) —
    # a real "well yes, but actually no" meme frame (pirate gesture): "would love
    # a magic pill — but no". Compact, does not compete with the mega-"03".
    _cover_pirate_meme(s, 0.75, 6.06)
    speaker_notes(s, load_notes("s02"))


def _cover_pirate_meme(s, x, y):
    """Small real internet-meme motif for the cover: pirate gesture "well yes, but
    actually no" (imgflip; English caption cropped out) + an English caption beside it
    — a hint at the lecture thesis: the "magic pill" does not exist. Compact:
    bottom-left corner, does not compete with the mega-"03". Attribution — attribution.md only."""
    mw = 2.10
    mh = mw * (755 / 1600)      # preserve crop proportions
    ocean_box(s, x - 0.08, y - 0.08, mw + 0.16, mh + 0.16)
    add_image(s, WEB / "cover-pirate-crop.png", x, y, mw, mh)
    text_box(s, x + mw + 0.28, y - 0.06, 5.4, mh + 0.16,
             "the “magic pill” does not exist",
             size=14, italic=True, bold=True, color=MID,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)


def build_s02a(p):
    """NEW (#212) lecture-map — 6 horizontal section cards (L1/L2 pattern).
    A separate contents slide after the cover. Shows the lecture route:
    Sections 0–5 with one line of meaning for each."""
    s = blank(p)
    slide_title(s, "Lecture route — six sections.", size=27)
    # #185/#315: the "Agents" section is no longer highlighted; the gold accent stays on
    # the lecture's load-bearing line (a marker before the subtitle), not on a section.
    filled_rect(s, 0.55, 1.26, 0.06, 0.34, GOLD)
    text_box(s, 0.74, 1.22, 12.05, 0.42,
             "One load-bearing line: picking the architecture for the task — and when the right answer is “not AI”.",
             size=15, italic=True, color=MID)
    cards = [
        ("0", "Opening", "framing the problem: where reliability comes from", "gavel", MID),
        ("1", "Prompt and its limits", "what one call can do and where its ceiling is", "message-circle", MID),
        ("2", "RAG", "external knowledge into context — and where it quietly breaks", "database", MID),
        ("3", "Fine-tuning", "changing weights for behavior, not for knowledge", "sliders-horizontal", MID),
        ("4", "Agents", "loop, equipment, memory, security — and failures", "bot", MID),
        ("5", "Framework", "ladder + checklist: how to choose fast and with grounds", "list-checks", MID),
    ]
    # 2 rows × 3 cols
    cw, chh = 3.95, 2.30
    gapx, gapy = 0.20, 0.28
    x0, y0 = 0.55, 1.90
    for i, (num, title, sub, ic, col) in enumerate(cards):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gapx)
        y = y0 + r * (chh + gapy)
        isgold = (col == GOLD)
        if isgold:
            ocean_box(s, x, y, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, y, cw, chh)
        circle(s, x + 0.24, y + 0.24, 0.56, col)
        text_box(s, x + 0.24, y + 0.24, 0.56, 0.56, num,
                 size=22, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        icon(s, ic, x + cw - 0.74, y + 0.26, 0.50, "gold" if isgold else "mid")
        text_box(s, x + 0.24, y + 0.98, cw - 0.48, 0.60, title,
                 size=16, bold=True, color=DEEP, line_spacing=1.05)
        text_box(s, x + 0.24, y + 1.58, cw - 0.48, 0.62, sub,
                 size=12.5, color=SLATE, line_spacing=1.14)
    speaker_notes(s, load_notes("s02a"))


def build_s03(p):
    """recap (§0) — a light reminder of two concepts from Lecture 2 that the
    whole lecture leans on: a single call (single-shot) and semantic search
    on embeddings. v6 (#185/#316): the former heavy "4 wrappers" schema removed —
    the wrappers are opened up section by section, here only the support from L2."""
    s = blank(p)
    slide_title(s, "From Lecture 2 we take two concepts as given.", size=27)
    text_box(s, 0.55, 1.22, 12.25, 0.42,
             "We are not re-explaining — just recalling, so we can lean on them further. Everything else we build up along the way.",
             size=15, italic=True, color=MID, line_spacing=1.15)
    # two large recap boxes
    by, bh = 2.05, 3.35
    bw = 6.05
    # LEFT — single-shot
    ocean_box(s, 0.55, by, bw, bh)
    icon(s, "message-circle", 0.85, by + 0.30, 0.62, "mid")
    text_box(s, 1.70, by + 0.30, bw - 1.10, 0.60, "Single call (single-shot)",
             size=19, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, 0.85, by + 1.20, bw - 0.60, 1.05,
             "One pass of the model: gave a prompt → got an answer. No memory between calls, no reaching outward.",
             size=15, color=DEEP, line_spacing=1.20)
    filled_rect(s, 0.85, by + 2.50, bw - 0.60, 0.66, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.5, radius=True, radius_adj=0.14)
    text_box(s, 1.05, by + 2.53, bw - 1.00, 0.60,
             "The model knows only what is in the prompt and in the weights",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    # RIGHT — semantic search / embeddings
    rx = 6.75
    ocean_box(s, rx, by, bw, bh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "database", rx + 0.30, by + 0.30, 0.62, "teal")
    text_box(s, rx + 1.15, by + 0.30, bw - 1.40, 0.60, "Semantic search on embeddings",
             size=19, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, rx + 0.30, by + 1.20, bw - 0.60, 1.05,
             "Text → a vector of meaning; fragments close in meaning are close vectors. Search “by meaning”, not by exact word.",
             size=15, color=DEEP, line_spacing=1.20)
    filled_rect(s, rx + 0.30, by + 2.50, bw - 0.60, 0.66, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.14)
    text_box(s, rx + 0.50, by + 2.53, bw - 1.00, 0.60,
             "RAG stands on this (Section 2)",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    gold_callout(s, 0.55, 5.72, 12.25, 0.88,
                 "Both concepts are the foundation from Lecture 2. Everything we build today (RAG, tools, fine-tuning, agents) leans on them.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """assertion_visual — central question + 6-step ladder."""
    s = blank(p)
    text_box(s, 0.55, 0.38, 12.25, 0.42, "THE CENTRAL QUESTION OF THE LECTURE",
             size=14, bold=True, color=TEAL)
    qx, qy, qw, qh = 0.55, 0.85, 12.25, 1.30
    ocean_box(s, qx, qy, qw, qh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, qx + 0.4, qy + 0.18, qw - 0.8, qh - 0.36,
             "I have a task and access to an LLM. Which architecture do I pick — and when is the right answer “not AI”?",
             size=23, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.12)
    # Ladder — 6 steps bottom-up: idx 0 = step 1 (BOTTOM, gold), idx 5 = step 6 (top).
    steps = [
        ("1", "Plain code (no AI)", "the reference point", GOLD, True),
        ("2", "One LLM call", "prompt; + CoT (step by step), + examples in the prompt (few-shot)", MID, False),
        ("3", "RAG / context engineering", "retrieval-augmented generation", LIGHT, False),
        ("4", "Workflow", "predefined paths", LIGHT, False),
        ("5", "Agent", "loop: plan → act → check → iterate", LIGHT, False),
        ("6", "Multi-agent", "several coordinated agents", LIGHT, False),
    ]
    n = len(steps)
    # #215/#216: the ladder's visual elements enlarged — bigger rungs
    # (step_h 0.66→0.74), bigger number circles (0.36→0.46) and sub-labels
    # (10.5→12), plus the direction arrow and its labels.
    step_h = 0.74
    vgap = 0.075
    bottom_edge = 6.85  # bottom of step 1
    sx = 0.55
    for i, (num, label, sub, col, isgold) in enumerate(steps):
        # i=0 -> bottom rung; each higher step sits above + indented right
        y = bottom_edge - step_h - i * (step_h + vgap)
        indent = i * 0.22
        w = 7.95 - indent
        x = sx + indent
        if isgold:
            ocean_box(s, x, y, w, step_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.5)
        else:
            ocean_box(s, x, y, w, step_h)
        circle(s, x + 0.16, y + (step_h - 0.46) / 2, 0.46, col)
        text_box(s, x + 0.16, y + (step_h - 0.46) / 2, 0.46, 0.46, num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.74, y + 0.11, w - 0.90, 0.30, label,
                 size=16, bold=True, color=DEEP)
        text_box(s, x + 0.74, y + 0.43, w - 0.90, 0.26, sub,
                 size=12, italic=True, color=SLATE)
    # PA-1 (owner-approved) + #216: direction-of-scale strip alongside ladder —
    # enlarged "more complex ↑" / "simpler ↓" + thicker arrow.
    text_box(s, 7.98, 2.28, 1.22, 0.34, "more complex ↑", size=14, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    up_arrow(s, 8.40, 2.72, 0.40, 3.86, fill=COVER_OUTLINE)
    text_box(s, 7.98, 6.62, 1.22, 0.34, "simpler ↓", size=14, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    # rule on the right
    ocean_box(s, 9.05, 2.55, 3.75, 3.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "milestone", 9.32, 2.85, 0.62, "gold")
    text_box(s, 9.32, 3.70, 3.25, 2.65,
             "Climb to the next step — only when the task demands something the current one does not cover.",
             size=16, bold=True, color=DEEP, line_spacing=1.24)
    footer(s, "The ladder is a map of the lecture, not a demand to grasp everything now. We work through each step separately.")
    speaker_notes(s, load_notes("s04"))


def build_s05(p):
    """assertion_visual — default = one call; cost of each climb."""
    s = blank(p)
    slide_title(s, "By default — one call with a good prompt.", size=27)
    # #219-context: explicit knowledge boundary of the model on the visible layer (§1.1)
    filled_rect(s, 0.55, 1.18, 12.25, 0.66, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.5, radius=True, radius_adj=0.10)
    text_box(s, 0.78, 1.24, 11.8, 0.56,
             "The model knows only what is in the prompt (plus what was in the weights at training). It has nowhere else to draw from.",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.12)
    # Left — anchor block
    lx, ly, lw, lh = 0.55, 2.02, 6.15, 3.62
    ocean_box(s, lx, ly, lw, lh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "target", lx + 0.28, ly + 0.24, 0.58, "teal")
    text_box(s, lx + 0.28, ly + 0.94, lw - 0.56, 0.6,
             "One LLM call\nwith a good prompt",
             size=18, bold=True, color=DEEP, line_spacing=1.1)
    bullets = [
        "lowest cost (a single pass)",
        "lowest latency (no extra round-trips)",
        "highest predictability (no loops, no retrieval that quietly degrades)",
    ]
    by = ly + 1.90
    for b in bullets:
        circle(s, lx + 0.30, by + 0.07, 0.12, TEAL)
        text_box(s, lx + 0.56, by, lw - 0.84, 0.62, b,
                 size=13, color=DEEP, line_spacing=1.12)
        by += 0.56
    # Right — a real "Gru's Plan" internet meme (imgflip, #185; replacing
    # Expanding Brain — owner: "the brain meme was already used"): escalating architecture
    # "just in case" with an absurd punchline in the 4th panel — the task was
    # three lines of plain code. English captions baked into the template.
    _s05_overengineering_meme(s, 6.95, 2.02, 5.85, 3.62)
    gold_callout(s, 0.55, 5.92, 12.25, 0.92,
                 "Do not complicate the architecture without a reason expressed in the task requirements. This is a distribution of the burden of proof, not primitivism.",
                 size=15)
    speaker_notes(s, load_notes("s05"))


def _s05_overengineering_meme(s, x, y, w, h):
    """A real "Gru's Plan" internet meme (imgflip blank + English captions):
    escalating architecture "just in case" with an absurd punchline in the 4th
    panel. Heading "over-engineering" + the meme itself in an Ocean frame. The meme
    carries the thesis "do not complicate without a reason" without on-slide attribution (attribution.md)."""
    ocean_box(s, x, y, w, h)
    chip(s, x + 0.28, y + 0.22, w - 0.56, 0.46, "“JUST IN CASE” — ESCALATION",
         fill=LIGHT, color=WHITE, size=13.5)
    # Gru's Plan (composite 700x449, landscape 1.56) fit to the box width
    from PIL import Image as _I
    _im = _I.open(WEB / "s05-gru-en.png"); _r = _im.size[0] / _im.size[1]; _im.close()
    img_w = w - 0.56
    img_h = img_w / _r
    if img_h > h - 0.86:
        img_h = h - 0.86
        img_w = img_h * _r
    img_x = x + (w - img_w) / 2
    add_image(s, WEB / "s05-gru-en.png", img_x, y + 0.82, img_w, img_h)


def build_s05a(p):
    """NEW (§1.2) — roles in the prompt: the myth "persona = accuracy" refuted.
    Zheng et al. 2024 EMNLP + arXiv:2605.29420. Part of the failure/judgment
    content — refuting the "magic pill"."""
    s = blank(p)
    slide_title(s, "A role in the prompt tunes tone — not accuracy.", size=26)
    text_box(s, 0.55, 1.16, 12.25, 0.42,
             "“You are an experienced lawyer” shifts the model's attention toward the text of such a role — but that is about style, not facts.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # #185 re-layout: minimum text + meme. Left third — myth + one experiment
    # fact; center — a large "Change my mind" meme; right — what the role
    # actually does. Gold takeaway as a wide strip under all three.
    ly = 1.78
    mx, mw = 5.20, 4.05
    mem_h = mw / 1.335                 # ≈3.03
    col_h = mem_h + 0.10              # cards to the meme height (mass balance)
    # LEFT — myth + one result (compressed)
    lx, lw = 0.55, 4.35
    ocean_box(s, lx, ly, lw, col_h)
    text_box(s, lx + 0.28, ly + 0.22, lw - 0.56, 0.36, "A common myth",
             size=15, bold=True, color=LIGHT)
    text_box(s, lx + 0.28, ly + 0.62, lw - 0.56, 0.80,
             "“An expert role in the prompt — and the model is more accurate on facts”",
             size=15, italic=True, color=SLATE, line_spacing=1.18)
    connector(s, lx + 0.28, ly + 1.58, lx + lw - 0.28, ly + 1.58, LIGHT, 1.0)
    text_box(s, lx + 0.28, ly + 1.74, lw - 0.56, 0.36, "Experiment",
             size=15, bold=True, color=MID)
    text_box(s, lx + 0.28, ly + 2.14, lw - 0.56, 0.90,
             "162 personas, 2410 questions —\npersonas did NOT raise accuracy.",
             size=15.5, bold=True, color=DEEP, line_spacing=1.26)
    # CENTER — a large "Change my mind" meme (its own full slot)
    my = ly
    ocean_box(s, mx - 0.14, my - 0.10, mw + 0.28, mem_h + 0.20)
    add_image(s, WEB / "s05a-changemymind-en.png", mx, my, mw, mem_h)
    # RIGHT — what the role actually does (height to the meme)
    rx, rw = 9.55, 3.25
    ocean_box(s, rx, ly, rw, col_h, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "user-round", rx + 0.26, ly + 0.24, 0.48, "teal")
    text_box(s, rx + 0.86, ly + 0.24, rw - 1.05, 0.48, "A role affects:",
             size=15, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.26, ly + 0.94, rw - 0.52, 1.20,
             "the tone and depth of the delivery — how formal and detailed, not the truth of a fact.",
             size=14.5, color=DEEP, line_spacing=1.24)
    # gold takeaway as a wide strip under the three columns (minimum text)
    gy = ly + col_h + 0.16
    gold_callout(s, 0.55, gy, 12.25, 6.90 - gy,
                 "Need accuracy? The tool is not a role, but context and RAG (a verifiable source). Another "
                 "“magic pill” point refuted by measurement.",
                 size=15)
    footer(s, "Zheng, Pei, Logeswaran, Lee, Jurgens · Findings of EMNLP 2024 (arXiv:2311.10054) + arXiv:2605.29420 (2026).")
    speaker_notes(s, load_notes("s05a"))


def build_s05b(p):
    """NEW (§1.3) — prompt structure: delimiters + separating
    instruction/context/data. A parallel with structured output (§4.1):
    input vs output — the same principle."""
    s = blank(p)
    slide_title(s, "Prompt structure: separate instruction, context, data.",
                size=23, h=1.30, line_spacing=1.08)
    text_box(s, 0.55, 1.58, 12.25, 0.42,
             "A flat prompt forces the model to guess where the instruction ends and the data begins. An explicit boundary removes that ambiguity.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # left — 3 stacked labelled parts (structured input)
    lx, ly, lw = 0.55, 2.14, 5.85
    parts = [
        ("braces", "Instruction", "what to do", MID),
        ("book-open", "Context", "on what basis", LIGHT),
        ("database", "Data", "what exactly to work with", TEAL),
    ]
    ph = 1.16
    pgap = 0.16
    py = ly
    for ic, t, sub, col in parts:
        ocean_box(s, lx, py, lw, ph)
        filled_rect(s, lx + 0.22, py + 0.22, 0.72, ph - 0.44, col, radius=True,
                    radius_adj=0.16)
        icon(s, ic, lx + 0.34, py + ph / 2 - 0.24, 0.48, "white")
        text_box(s, lx + 1.14, py + 0.20, lw - 1.35, 0.42, t,
                 size=17, bold=True, color=DEEP)
        text_box(s, lx + 1.14, py + 0.62, lw - 1.35, 0.40, sub,
                 size=13, italic=True, color=SLATE)
        py += ph + pgap
    # right — delimiters kinds + structured-output parallel
    rx, rw = 7.15, 5.65
    ocean_box(s, rx, ly, rw, 1.95)
    text_box(s, rx + 0.28, ly + 0.18, rw - 0.56, 0.36, "Delimiters",
             size=15, bold=True, color=MID)
    for i, d in enumerate([
        "XML tags:  <instruction>…</instruction>",
        "Markdown headings:  ## Task · ## Data",
        "Triple quotes / backticks around the data",
    ]):
        circle(s, rx + 0.30, ly + 0.62 + i * 0.40 + 0.05, 0.12, LIGHT)
        text_box(s, rx + 0.54, ly + 0.62 + i * 0.40, rw - 0.82, 0.38, d,
                 size=13, color=DEEP, font=FONT_MONO if i < 2 else FONT_BODY,
                 line_spacing=1.05)
    ocean_box(s, rx, ly + 2.10, rw, 1.70, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.28, ly + 2.24, rw - 0.56, 1.44,
             "The same principle as structured output (Section 4): there a schema sets the structure of the OUTPUT, here delimiters set the structure of the INPUT. A clear structure instead of “infer the form by meaning”.",
             size=13.5, color=DEEP, line_spacing=1.20)
    gold_callout(s, 0.55, 5.98, 12.25, 0.86,
                 "A model that has been shown where the instruction ends and the data begins less often takes a fragment of data for a new command — the same confusion underlies prompt injection (Section 4).",
                 size=13.5)
    speaker_notes(s, load_notes("s05b"))


def build_s06(p):
    """case_study — CoT worked example + faithfulness limit MERGED (§1.4+§1.5)."""
    s = blank(p)
    slide_title(s, "Chain-of-thought helps — but you cannot audit it.", size=24)
    # TOP BAND — CoT worked example (compact before/after)
    cy, ch = 1.16, 2.02
    cw = 6.05
    ocean_box(s, 0.55, cy, cw, ch)
    text_box(s, 0.83, cy + 0.14, cw - 0.56, 0.34, "Without CoT",
             size=15, bold=True, color=LIGHT)
    text_box(s, 0.83, cy + 0.52, cw - 0.56, 0.85,
             "“There were 23 apples, 7 spoiled, then 2 crates of 6 were bought. How many good ones?”",
             size=13, color=DEEP, line_spacing=1.14)
    text_box(s, 0.83, cy + 1.42, cw - 0.56, 0.5,
             "→ a plausible but wrong number",
             size=13, bold=True, color=SLATE)
    rx0 = 6.75
    ocean_box(s, rx0, cy, cw, ch, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx0 + 0.28, cy + 0.14, cw - 0.56, 0.34, "With CoT (“reason step by step”)",
             size=15, bold=True, color=TEAL)
    text_runs(s, rx0 + 0.30, cy + 0.56, cw - 0.6, 0.9, [
        {"text": "23 − 7 = 16    ·    2 × 6 = 12", "size": 15, "bold": True, "color": DEEP},
        {"text": "16 + 12 = 28", "size": 22, "bold": True, "color": GOLD,
         "newpara": True, "space_before": 6},
    ], line_spacing=1.1)
    text_box(s, rx0 + 0.30, cy + 1.56, cw - 0.6, 0.36, "→ correct",
             size=13, bold=True, color=TEAL)
    # thin note between bands (single line)
    text_box(s, 0.55, cy + ch + 0.06, 12.25, 0.28,
             "Technically this is still one call — CoT is a tool for a class of tasks (a chain of steps), not a global switch.",
             size=12, italic=True, color=MID)
    # BOTTOM BAND — faithfulness limit (qualitative; numbers — on the next slide)
    fy = 3.70
    text_box(s, 0.55, fy, 12.25, 0.34,
             "But the stated reasoning need not reflect the real cause of the answer (low explanation faithfulness):",
             size=14, bold=True, color=DEEP)
    st_y = fy + 0.42
    # #185/#319: the right text block “Control over self-explanation…” is
    # REMOVED entirely — the point is carried by the real internet meme
    # “Distracted boyfriend” (imgflip): the model got distracted by a nice
    # out-loud explanation and missed the real cause of the answer. Captions baked into the template.
    mh = 2.64
    mw = mh * (1200 / 800)
    mx = 0.55
    ocean_box(s, mx - 0.12, st_y - 0.10, mw + 0.24, mh + 0.20)
    add_image(s, WEB / "s06-distracted-en.png", mx, st_y, mw, mh)
    # short takeaway on the right (replaces the removed text block); gold accent on the verb
    rx = mx + mw + 0.44
    rw = 13.33 - rx - 0.55
    gh = 1.72
    gy = st_y + (mh - gh) / 2
    filled_rect(s, rx, gy, rw, gh, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_runs(s, rx + 0.34, gy + 0.14, rw - 0.68, gh - 0.28, [
        {"text": "Check the result", "size": 21, "bold": True, "color": DEEP},
        {"text": ", not the model’s self-explanation.", "size": 21, "bold": True,
         "color": DEEP},
        {"text": "Verify facts against an external source — a plausible reasoning text confirms nothing.",
         "size": 15, "bold": False, "color": DEEP, "newpara": True,
         "space_before": 8},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
    footer(s, "Anthropic, April 2025 · faithfulness drops on hard tasks · re-verify before the lecture day.")
    speaker_notes(s, load_notes("s06"))


def build_s08(p):
    """assertion_visual — context engineering + context rot curve (c08)."""
    s = blank(p)
    slide_title(s, "Context engineering: the minimum of high-signal.", size=26)
    text_box(s, 0.55, 1.20, 12.25, 0.55,
             "Prompt engineering is one instruction. Context engineering is curating the whole set of tokens the model sees at inference.",
             size=14, italic=True, color=MID, line_spacing=1.15)
    # left — curve chart
    cx, cyy, cw, chh = 0.55, 1.95, 7.05, 3.55
    ocean_box(s, cx, cyy, cw, chh)
    add_image(s, CHARTS / "c08-context-rot.png", cx + 0.18, cyy + 0.16,
              cw - 0.36, chh - 0.32)
    text_box(s, cx + 0.18, cyy + chh + 0.02, cw - 0.36, 0.42,
             "context rot = the same “lost in the middle” from L2 — a new term, not a new entity",
             size=11.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # right — criterion
    rx, rw = 7.85, 4.95
    ocean_box(s, rx, 1.95, rw, 3.55, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.28, 2.16, rw - 0.56, 0.45, "When NOT RAG (point 1)",
             size=17, bold=True, color=TEAL)
    text_box(s, rx + 0.28, 2.70, rw - 0.56, 1.55,
             "a small, stable corpus that fits in the window → full context + prefix caching, not RAG infrastructure",
             size=15, color=DEEP, line_spacing=1.22)
    icon(s, "circle-slash", rx + 0.28, 4.45, 0.78, "teal")
    text_box(s, rx + 1.20, 4.55, rw - 1.45, 0.85,
             "RAG here would add fragility without a payoff",
             size=13.5, bold=True, color=DEEP, line_spacing=1.12,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.95, 12.25, 0.88,
                 "“Find the smallest set of high-signal tokens that maximizes the probability of the desired outcome” — this is an engineering requirement, not aesthetics.",
                 size=14.5)
    speaker_notes(s, load_notes("s08"))


def build_s08a(p):
    """NEW (§1.8) — cheat sheet “how to build a prompt”, 8 items. A compact
    analogue of the §5.3 checklist at the single-prompt level."""
    s = blank(p)
    slide_title(s, "Cheat sheet: how to build a prompt.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.40,
             "The minimum below which a prompt systematically underperforms. Apply it to any prompt before the first run.",
             size=14, italic=True, color=MID)
    items = [
        ("Role", "if you need tone/register — and NOT as a promise of accuracy"),
        ("Task", "a concrete verifiable action, not a vague wish"),
        ("Context", "the minimum needed, not “everything that might come in handy”"),
        ("Output format", "stated explicitly if the answer is machine-processed"),
        ("Delimiters", "if there is more than one kind of content (instruction/data)"),
        ("Examples (few-shot)", "only if the format is not obvious from one instruction"),
        ("CoT", "only if the task needs multi-step reasoning"),
        ("Length", "no longer than needed — extra tokens “drown” in the context"),
    ]
    cw, chh = 3.95, 1.28
    gapx, gapy = 0.20, 0.20
    x0, y0 = 0.55, 1.78
    for i, (t, sub) in enumerate(items):
        r, c = divmod(i, 2)
        # 2 cols × 4 rows
        x = x0 + c * (cw + gapx) if False else 0.0
        # 4 cols × 2 rows layout
        col = i % 4
        row = i // 4
        cw2 = 3.02
        x = 0.55 + col * (cw2 + 0.13)
        y = y0 + row * (chh + gapy)
        isgold = (i == 0)  # role = the key item (the accuracy myth)
        if isgold:
            ocean_box(s, x, y, cw2, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, y, cw2, chh)
        circle(s, x + 0.18, y + 0.18, 0.34, GOLD if isgold else MID)
        text_box(s, x + 0.18, y + 0.18, 0.34, 0.34, str(i + 1),
                 size=13, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.62, y + 0.16, cw2 - 0.78, 0.40, t,
                 size=14, bold=True, color=DEEP)
        text_box(s, x + 0.20, y + 0.58, cw2 - 0.40, 0.62, sub,
                 size=11, color=SLATE, line_spacing=1.10)
    gold_callout(s, 0.55, 5.20, 12.25, 1.05,
                 "A compact form of the whole section: role ≠ accuracy, structure helps separate inputs, CoT is a targeted tool, context is minimal. For large architectural decisions (RAG / FT / agent) — the eight-step checklist of Section 5.",
                 size=14)
    speaker_notes(s, load_notes("s08a"))


def build_s09(p):
    """section_divider — Section 2 RAG."""
    build_section_divider(
        p, 2, "Section 2", "RAG: retrieval-augmented generation",
        "Retrieve the relevant → put it in context → answer grounded in a source",
        "s09",
        image_src=WEB / "div-r2-library.jpg",
        tag="external knowledge · 3 cases · 1 failure")


def build_s10(p):
    """schema_pipeline — RAG 3-stage horizontal pipeline (RIGHT_ARROW)."""
    s = blank(p)
    slide_title(s, "The RAG principle — three steps.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "RAG = indexing → retrieval (semantic search from L2) → grounded generation; “I don't know” is a correct answer.",
             size=14.5, italic=True, color=MID)
    # 3 stage boxes with RIGHT_ARROW between
    sy, sh = 1.85, 3.05
    bw = 3.55
    gap_arrow = 0.55
    x0 = 0.55
    stages = [
        ("1", "Indexing", "ahead of time, offline", "database",
         "corpus → chunks → embed each → vector store", MID, False),
        ("2", "Retrieval", "at query time", "route",
         "question → embedding → k nearest fragments\n\n= the same semantic search from L2 — not re-explained", MID, False),
        ("3", "Generation", "grounded in a source (grounding)", "check-check",
         "fragments + question → answer with a source reference", TEAL, True),
    ]
    x = x0
    for i, (num, title, tag, ic, body, col, isteal) in enumerate(stages):
        if isteal:
            ocean_box(s, x, sy, bw, sh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
        else:
            ocean_box(s, x, sy, bw, sh)
        circle(s, x + 0.26, sy + 0.26, 0.46, col)
        text_box(s, x + 0.26, sy + 0.26, 0.46, 0.46, num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.86, sy + 0.26, bw - 1.05, 0.42, title,
                 size=17, bold=True, color=DEEP)
        text_box(s, x + 0.86, sy + 0.66, bw - 1.05, 0.30, tag,
                 size=11.5, italic=True, color=LIGHT)
        icon(s, ic, x + bw - 0.78, sy + 0.28, 0.50,
             "teal" if isteal else "mid")
        text_box(s, x + 0.28, sy + 1.18, bw - 0.56, sh - 1.36, body,
                 size=13.5, color=DEEP, line_spacing=1.22)
        if i < 2:
            right_arrow(s, x + bw + 0.06, sy + sh / 2 - 0.30, gap_arrow - 0.12, 0.60,
                        fill=LIGHT)
        x += bw + gap_arrow
    gold_callout(s, 0.55, 5.10, 12.25, 0.90,
                 "“I don't know” / “see source X” is a correct answer from a RAG system. A plausible answer on an irrelevant retrieval is a defect, not “at least something”.",
                 size=14)
    # RAG-2026: current default stack (glossary-locked terms on the visible layer)
    filled_rect(s, 0.55, 6.12, 12.25, 0.62, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_runs(s, 0.78, 6.20, 11.8, 0.48, [
        {"text": "RAG-2026: ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": "agentic RAG by default · hybrid search (BM25 + dense vectors) · reranker · miss cascade ", "size": 12.5, "color": DEEP},
        {"text": "5.7% → 1.9%", "size": 12.5, "bold": True, "color": DEEP},
        {"text": " (Contextual Retrieval).", "size": 12.5, "color": DEEP},
    ], line_spacing=1.1)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """assertion_visual — 4 conjunction cards -> RAG."""
    s = blank(p)
    slide_title(s, "When RAG is the right choice.", size=27)
    # #221/#222: the jargon “conjunction/disjunction” is removed; the §2.2
    # wording — a strong signal on the features + no blockers from §2.3.
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "RAG is justified when the signal on the features below is strong — and there are no blockers from the next slide, “when NOT RAG”.",
             size=14.5, italic=True, color=MID)
    cards = [
        ("Large / growing", "does not fit the window whole, or it is costly to put the entire corpus in every request"),
        ("Changing", "documents, prices, regulations update more often than model versions ship"),
        ("Freshness + provenance", "the answer relies on a verifiable source; you can show where the fact came from"),
        ("Private base", "the company's knowledge is not in the weights of a public model"),
    ]
    cw, chh = 2.90, 2.55
    cy = 1.78
    x = 0.55
    for i, (t, sub) in enumerate(cards):
        ocean_box(s, x, cy, cw, chh)
        text_box(s, x + 0.20, cy + 0.22, cw - 0.40, 0.75, t,
                 size=15.5, bold=True, color=MID, line_spacing=1.05)
        text_box(s, x + 0.20, cy + 1.00, cw - 0.40, chh - 1.2, sub,
                 size=12, color=DEEP, line_spacing=1.16)
        x += cw + 0.20
    # -> RAG result + worked example
    ocean_box(s, 0.55, 4.55, 7.55, 1.80, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, 0.85, 4.72, 7.1, 0.4, "The features reinforce each other → RAG",
             size=17, bold=True, color=DEEP)
    text_box(s, 0.85, 5.16, 7.1, 1.10,
             "A corporate base of thousands of regulations, updated weekly, an answer with a mandatory reference to the clause: all features converge, no simpler mechanism closes them jointly → a textbook RAG profile.",
             size=13, color=DEEP, line_spacing=1.16)
    ocean_box(s, 8.30, 4.55, 4.50, 1.80, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "check-check", 8.55, 4.74, 0.44, "teal")
    text_box(s, 8.55, 5.24, 4.05, 1.00,
             "One feature is a reason to take a look, not to build automatically: check the blockers on the next slide.",
             size=12.5, bold=True, color=DEEP, line_spacing=1.15)
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """schema_matrix-ish — 3 criteria columns «when NOT RAG».
    P1 fix (issue #157 review): s11/s12 shared identical «N cards + summary
    plaque» skeleton back-to-back — visually merged into one stretched
    slide. Differentiated here via base card palette (TEAL instead of
    primary LIGHT/SURFACE — «caution / exclusion» register vs s11's
    primary-blue «inclusion» cards) + numbered gold badge per card (s22b
    slot-badge pattern) so the two decks read as distinct compositions at
    a glance, without touching content/copy."""
    s = blank(p)
    slide_title(s, "When RAG is NOT the right choice.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "Knowing when RAG is not needed is worth more: it is a fashionable architecture, put where it does harm.",
             size=14.5, italic=True, color=MID)
    # #223/#224: the observability criterion is removed (it lives on s13, the RAG
    # failure at scale, not duplicated here). In its place — a new criterion
    # “data available live via API/MCP” (§2.3, forward-callback to §4.1).
    cols = [
        ("circle-slash", "The corpus fits in the window",
         "a rough guide — under ~200k tokens, rarely changes",
         "→ full context + prefix caching, not RAG infrastructure",
         False),
        ("key", "A fixed policy / value",
         "a fare, a price, a clause of a regulation, a rule",
         "→ a deterministic lookup table / a static page",
         False),
        ("cable", "Data available in real time via API / MCP",
         "in an internal service, a database, another system's search",
         "→ call the tool directly; a RAG index on top is an extra layer, more fragile and more stale",
         True),
    ]
    # #185: the cards a bit lower → opening a row for the real “Roll Safe” meme.
    cw, chh = 4.00, 2.62
    cy = 1.78
    x = 0.55
    for i, (nm, t, tag, alt, isgold) in enumerate(cols):
        if isgold:
            ocean_box(s, x, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, cy, cw, chh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=1.5)
        # numbered badge top-left (s22b slot-badge pattern) — distinct
        # silhouette from s11's plain unnumbered cards
        circle(s, x + 0.20, cy + 0.18, 0.32, GOLD if isgold else TEAL)
        text_box(s, x + 0.20, cy + 0.18, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        icon(s, nm, x + cw - 0.76, cy + 0.20, 0.52, "gold" if isgold else "teal")
        text_box(s, x + 0.66, cy + 0.24, cw - 1.55, 0.78, t,
                 size=15, bold=True, color=(DEEP if isgold else TEAL),
                 line_spacing=1.08)
        text_box(s, x + 0.24, cy + 1.15, cw - 0.48, 0.40, tag,
                 size=12, italic=True, color=LIGHT)
        text_box(s, x + 0.24, cy + 1.58, cw - 0.48, chh - 1.72, alt,
                 size=13, bold=True, color=DEEP, line_spacing=1.14)
        x += cw + 0.32
    # #185: the real internet meme “Roll Safe” — “no RAG needed if the corpus
    # fits in the context” (the first of three blockers, the most common anti-pattern).
    rmw, rmh = 2.15, 1.21
    rmx, rmy = 0.55, 4.50
    filled_rect(s, rmx - 0.05, rmy - 0.05, rmw + 0.10, rmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s12-rollsafe-en.png", rmx, rmy, rmw, rmh)
    text_box(s, rmx + rmw + 0.36, rmy + 0.02, 9.05, rmh,
             "A common anti-pattern — putting RAG infrastructure where the whole corpus comfortably fits in the context window and rarely changes.",
             size=13.5, bold=True, color=MID, line_spacing=1.24,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.78, 12.25, 1.02,
                 "RAG is redundant if ANY of the three holds: (a) the corpus fits in the window and is stable, (b) the task reduces to a fixed value, (c) the knowledge is already available directly and live via a tool. “A direct call returns data as of the request; a RAG index — as of the last indexing.”",
                 size=13.5)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """case_study — RAG fails at scale + Air Canada — architecture walkthrough (v2: decluttered
    — per-case triangle-alert icons, tighter visible text, mass-rebalanced
    right box so diagnosis+alternative fill evenly, no big internal gap)."""
    s = blank(p)
    slide_title(s, "The failure of RAG at scale.", size=27)
    gold_callout(s, 0.55, 1.10, 12.25, 0.95,
                 "“Returned something” ≠ “returned the right thing”. RAG has no “not found” signal — it always returns the k nearest, even irrelevant ones.",
                 size=15)
    # left — 3 failure cases, each with a triangle-alert anchor icon
    # PA-2 (owner-approved): cases distributed EVENLY across the full box
    # height with thin separators (was ~15-20% dead band at the bottom;
    # mass now matches the right Air Canada box).
    lx, lw = 0.55, 6.55
    box_y, box_h = 2.14, 3.74
    ocean_box(s, lx, box_y, lw, box_h)
    cases = [
        ("Legal-AI", "the “nearest” cases from another jurisdiction / an overturned precedent — the model relies on them as fact"),
        ("Medical-RAG", "mixed fragments of different patients — close in symptoms, but clinically cannot be combined"),
        ("Support bot", "worked on hundreds of articles; after growing to thousands, quality quietly degraded — nobody noticed"),
    ]
    cell_h = box_h / 3.0          # 3 equal cells fill the box top→bottom
    for ci, (nm, body) in enumerate(cases):
        cy0 = box_y + ci * cell_h
        icon(s, "triangle-alert", lx + 0.28, cy0 + 0.22, 0.40, "teal")
        text_box(s, lx + 0.84, cy0 + 0.22, lw - 1.10, 0.36, nm,
                 size=16, bold=True, color=MID)
        text_box(s, lx + 0.28, cy0 + 0.66, lw - 0.54, 0.62, body,
                 size=13, color=DEEP, line_spacing=1.20)
        if ci < 2:
            connector(s, lx + 0.30, box_y + (ci + 1) * cell_h,
                      lx + lw - 0.30, box_y + (ci + 1) * cell_h, LIGHT, 0.75)
    # right — Air Canada — architecture walkthrough (teal), evenly split: diagnosis / alternative
    rx, rw = 7.35, 5.45
    ocean_box(s, rx, 2.14, rw, 3.74, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.30, 2.36, rw - 0.60, 0.34, "Air Canada — architecture walkthrough",
             size=16, bold=True, color=TEAL)
    text_box(s, rx + 0.30, 2.78, rw - 0.60, 0.30, "Diagnosis",
             size=13, bold=True, color=DEEP)
    text_box(s, rx + 0.30, 3.10, rw - 0.60, 1.05,
             "generated plausible text was placed in a role that required a retrieved, verified fact — a failure of grounding",
             size=13, color=DEEP, line_spacing=1.20)
    connector(s, rx + 0.30, 4.22, rx + rw - 0.30, 4.22, TEAL, 1.0)
    text_box(s, rx + 0.30, 4.34, rw - 0.60, 0.30, "The right alternative",
             size=13, bold=True, color=DEEP)
    text_box(s, rx + 0.30, 4.66, rw - 0.60, 1.15,
             "a fixed policy → a lookup table / page; need a dialogue → RAG with strict grounding, a mandatory citation, an explicit “I don't know”, a human check",
             size=13, color=DEEP, line_spacing=1.20)
    footer(s, "Documented failure classes (Barnett et al. 2024; Air Canada — McCarthy Tétrault 2024). Cases are illustrative.")
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """assertion_visual (§3.1/§3.5) — distillation = a SEPARATE technique, NOT a kind of
    fine-tuning (#227 P0). Schema: teacher (fine-tuned) → distill → student
    (smaller, cheaper). Contrast “what fine-tuning does / what distillation does”.
    Criteria “what goes where” — on s17."""
    s = blank(p)
    slide_title(s, "Distillation is not a kind of fine-tuning — it is a separate technique.", size=25)
    text_box(s, 0.55, 1.12, 12.25, 0.44,
             "Fine-tuning changes the model's behavior. Distillation is compression: transferring the skills of a large model into a small one. They are often confused, but these are two taxonomically distinct operations that work together.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # pipeline: teacher (FT) -> distill -> student  (compressed to the left — the right
    # column is given to the real “Yoda” meme teacher→student)
    sy, sh = 1.90, 2.48
    ocean_box(s, 0.55, sy, 8.55, sh)
    bw, bh = 2.05, 1.50
    by = sy + 0.56
    t1x = 0.85
    arr1 = t1x + bw
    t2x = arr1 + 0.80
    arr2 = t2x + bw
    t3x = arr2 + 0.80
    # teacher — big, fine-tuned
    filled_rect(s, t1x, by, bw, bh, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.08)
    icon(s, "cpu", t1x + bw / 2 - 0.20, by + 0.12, 0.38, "teal")
    text_box(s, t1x + 0.08, by + 0.58, bw - 0.16, 0.30, "Teacher",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, t1x + 0.10, by + 0.90, bw - 0.20, 0.52,
             "large, fine-tuned for the task", size=10.5, italic=True,
             color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.04)
    text_box(s, arr1 - 0.32, by - 0.40, 1.25, 0.30, "distillation",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    right_arrow(s, arr1 + 0.06, by + bh / 2 - 0.17, 0.68, 0.34, fill=MID)
    # student — small, cheaper (gold anchor)
    filled_rect(s, t2x, by + 0.20, bw, bh - 0.40, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.08)
    icon(s, "cpu", t2x + bw / 2 - 0.17, by + 0.28, 0.32, "gold")
    text_box(s, t2x + 0.08, by + 0.66, bw - 0.16, 0.30, "Student",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, t2x + 0.10, by + 0.96, bw - 0.20, 0.44,
             "small, cheaper and faster", size=10.5,
             italic=True, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.04)
    text_box(s, arr2 - 0.32, by - 0.40, 1.25, 0.30, "to production →",
             size=11, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    right_arrow(s, arr2 + 0.06, by + bh / 2 - 0.17, 0.68, 0.34, fill=LIGHT)
    filled_rect(s, t3x, by, bw, bh, SURFACE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.08)
    icon(s, "target", t3x + bw / 2 - 0.20, by + 0.12, 0.38, "mid")
    text_box(s, t3x + 0.08, by + 0.58, bw - 0.16, 0.30, "In production",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, t3x + 0.10, by + 0.90, bw - 0.20, 0.52,
             "lower cost and latency", size=10.5,
             italic=True, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.04)
    # #185: the real internet meme “Yoda” — the teacher (a large fine-tuned
    # model) passes the skill to a small student (distillation). Right column,
    # full height of the schema + contrast plaques.
    ymx, ymy, ymw, ymh = 9.35, 1.90, 3.45, 4.08
    ocean_box(s, ymx, ymy, ymw, ymh)
    add_image(s, WEB / "s14-yoda-en.png", ymx + 0.14, ymy + 0.14,
              ymw - 0.28, ymh - 0.28)
    # contrast strip (narrowed to the left — the right column is given to the meme)
    cy, ch = 4.56, 1.42
    ocean_box(s, 0.55, cy, 4.20, ch)
    text_box(s, 0.78, cy + 0.14, 3.80, 0.34, "Fine-tuning answers:",
             size=13, bold=True, color=MID)
    text_box(s, 0.78, cy + 0.50, 3.80, 0.86,
             "“how the model behaves” — tone, format, policy. Changes the WEIGHTS for behavior.",
             size=12, color=DEEP, line_spacing=1.16)
    ocean_box(s, 4.90, cy, 4.20, ch, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, 5.13, cy + 0.14, 3.80, 0.34, "Distillation answers:",
             size=13, bold=True, color=TEAL)
    text_box(s, 5.13, cy + 0.50, 3.80, 0.86,
             "“how to make the model cheaper” — the same skill in a smaller model. This is compression, not a change of behavior.",
             size=12, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 6.14, 12.25, 0.78,
                 "What to do: do not write “distillation is fine-tuning”. In tandem they go like this: first fine-tune the teacher for the task, then distill into the student for cost. The criteria “what goes where” — on the next slide.",
                 size=13)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    """assertion_visual — PEFT vs full-FT: frozen base + adapters + 3 reasons."""
    s = blank(p)
    slide_title(s, "PEFT instead of full fine-tuning.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.78,
             "PEFT — the base weights are frozen, only a small set of adapters is trained. LoRA — low-rank adapter matrices; QLoRA — the same on top of a quantized model.",
             size=14, italic=True, color=MID, line_spacing=1.18)
    # left — schema: big frozen base + small adapters
    lx, ly, lw, lh = 0.55, 2.10, 4.95, 3.15
    ocean_box(s, lx, ly, lw, lh)
    filled_rect(s, lx + 0.55, ly + 0.55, lw - 1.1, 1.55, MID, radius=True,
                radius_adj=0.06)
    text_box(s, lx + 0.55, ly + 0.55, lw - 1.1, 1.55,
             "Base weights\n(frozen)", size=17, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    for k in range(3):
        ax = lx + 0.70 + k * 1.20
        filled_rect(s, ax, ly + 2.35, 0.95, 0.55, GOLD, radius=True,
                    radius_adj=0.18)
        text_box(s, ax, ly + 2.35, 0.95, 0.55, "LoRA",
                 size=11, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx, ly + lh - 0.02, lw, 0.3,
             "adapters — megabytes vs gigabytes", size=11, italic=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    # right — 3 reasons (more compact — freeing the right column for the meme)
    rx, rw = 5.80, 3.55
    reasons = [
        ("1. Cheaper and faster", "millions of parameters instead of billions; QLoRA — on a single GPU", False),
        ("2. Modularity", "adapters are megabytes vs gigabytes; one base — many specializations", False),
        ("3. ↓ Forgetting risk", "the base is frozen, physically not overwritten by the new signal", True),
    ]
    yy = 2.10
    for t, b, isgold in reasons:
        if isgold:
            ocean_box(s, rx, yy, rw, 1.02, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, rx, yy, rw, 0.98)
        text_box(s, rx + 0.22, yy + 0.10, rw - 0.44, 0.34, t,
                 size=14, bold=True, color=(DEEP if isgold else MID))
        text_box(s, rx + 0.22, yy + 0.44, rw - 0.44, (0.52 if isgold else 0.48), b,
                 size=12, color=DEEP, line_spacing=1.10)
        yy += (1.14 if isgold else 1.10)
    # #185: the real internet meme “Buff Doge vs Cheems” — PEFT (a strong,
    # calm choice) vs full fine-tuning (expensive, weak). Right column.
    dmx, dmy, dmw, dmh = 9.55, 2.10, 3.25, 3.15
    ocean_box(s, dmx, dmy, dmw, dmh)
    add_image(s, WEB / "s15-doge-en.png", dmx + 0.12, dmy + 0.12,
              dmw - 0.24, dmh - 0.24)
    # LoRA-adoption baseline (§3.2) with a MANDATORY caveat on the visible layer
    # (Baseline Mandate): the share among PEFT-tagged models, not among all FT.
    by2 = 5.50
    ocean_box(s, 0.55, by2, 5.95, 1.32, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_runs(s, 0.80, by2 + 0.16, 5.5, 0.55, [
        {"text": "98.4% ", "size": 30, "bold": True, "color": GOLD},
        {"text": "of PEFT-tagged models are LoRA", "size": 14, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.80, by2 + 0.72, 5.45, 0.52,
             "out of 20,834 cards on the Hugging Face Hub · caveat: the share among PEFT-tagged models, not among all fine-tuning",
             size=11, italic=True, color=SLATE, line_spacing=1.12)
    gold_callout(s, 6.70, by2, 6.10, 1.32,
                 "PEFT (LoRA/QLoRA) is almost always preferable to full fine-tuning: cheaper, more modular, ↓ forgetting risk. Full fine-tuning in 2026 — almost never.",
                 size=14)
    footer(s, "HF PEFT team, “Beyond LoRA?”, June 2026. The full spectrum of methods (SFT / DPO / RFT) — in the chapter.")
    speaker_notes(s, load_notes("s15"))


def build_s16(p):

    """case_study — catastrophic forgetting (diverging chart c16) + criterion."""
    s = blank(p)
    slide_title(s, "Failure: catastrophic forgetting.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.55,
             "Catastrophic forgetting — degradation of the model's general abilities as a result of narrow aggressive fine-tuning.",
             size=14, italic=True, color=MID, line_spacing=1.15)
    cx, cyy, cw, chh = 0.55, 1.92, 7.05, 3.55
    ocean_box(s, cx, cyy, cw, chh)
    add_image(s, CHARTS / "c16-forgetting.png", cx + 0.18, cyy + 0.16,
              cw - 0.36, chh - 0.32)
    text_box(s, cx + 0.18, cyy + chh + 0.02, cw - 0.36, 0.42,
             "worse as the model scales up — a larger model starts higher, so it has \"further to fall\"",
             size=11.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    rx, rw = 7.85, 4.95
    ocean_box(s, rx, 1.92, rw, 2.30, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.26, 2.10, rw - 0.52, 2.0,
             "No eval loop on general tasks + no dataset/weights versioning → you won't see the break before production + you can't roll back → this is not a \"risk\", it is a criterion for \"do NOT fine-tune\"",
             size=13.5, bold=True, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, 4.37, rw, 1.10, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.26, 4.52, rw - 0.52, 0.85,
             "Right way: PEFT (frozen weights — lower risk); for changing knowledge — RAG, not fine-tuning at all.",
             size=12.5, color=DEEP, line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Empirically observed under continual fine-tuning (Luo et al. 2023). Mechanisms — \"research suggests\" (preprint).")
    speaker_notes(s, load_notes("s16"))


def build_s18(p):
    """section_divider — Section 4 "Agents" (title WITHOUT "+security";
    security content lives inside the section, on s25)."""
    build_section_divider(
        p, 4, "Section 4", "Agents",
        "From a conversationalist in a chat window — to a component of a production system: the loop, the equipment, memory, tool access — and where all of it breaks.",
        "s18",
        image_src=WEB / "div-r4-robot-arm.jpg",
        tag="agent + equipment · 5 cases · 4 failures")


def build_s19(p):
    """MERGED (§4.1) — API mechanics (structured output / function calling /
    prompt caching) ON TOP + MCP (N×M→N+M, USB-C, adoption, trust shift)
    BELOW. Dense slide — compact cards, separate 5-Second check."""
    s = blank(p)
    slide_title(s, "The model becomes a system component: API + MCP.", size=24)
    # TOP — 3 API mechanism cards (compact)
    cards = [
        ("boxes", "Structured output", "output strictly by schema (JSON), not text to parse", "embeddable"),
        ("terminal", "Tool calling", "the model formulates \"call X\"; your code executes it, not the model", "active"),
        ("database", "Prompt caching", "don't recompute the unchanged prefix on every request", "economical"),
    ]
    cw, chh = 4.00, 2.02
    cy = 1.12
    x = 0.55
    for nm, t, body, tag in cards:
        ocean_box(s, x, cy, cw, chh)
        filled_rect(s, x + 0.22, cy + 0.22, 0.56, 0.56, MID, radius=True, radius_adj=0.18)
        icon(s, nm, x + 0.28, cy + 0.28, 0.44, "white")
        text_box(s, x + 0.90, cy + 0.24, cw - 1.05, 0.52, t,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x + 0.24, cy + 0.90, cw - 0.48, 0.72, body,
                 size=12.5, color=DEEP, line_spacing=1.14)
        chip(s, x + 0.24, cy + chh - 0.46, cw - 0.48, 0.34, tag,
             fill=TEAL, color=WHITE, size=12)
        x += cw + 0.13
    # BOTTOM — MCP
    my = 3.36
    ocean_box(s, 0.55, my, 6.55, 2.55)
    icon(s, "cable", 0.80, my + 0.20, 0.48, "mid")
    text_runs(s, 1.40, my + 0.22, 5.5, 0.5, [
        {"text": "MCP", "size": 18, "bold": True, "color": MID},
        {"text": "  — \"USB-C for LLM tools\"", "size": 14, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text_runs(s, 0.80, my + 0.82, 6.0, 0.5, [
        {"text": "N×M", "size": 20, "bold": True, "color": LIGHT},
        {"text": " incompatible integrations → ", "size": 13, "color": DEEP},
        {"text": "N+M", "size": 20, "bold": True, "color": TEAL},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.80, my + 1.36, 6.0, 0.32,
             "a tool once = an MCP server; a model once = an MCP client",
             size=11.5, italic=True, color=SLATE)
    tl = [("Anthropic", "11/2024", "logo-anthropic"),
          ("OpenAI", "03/2025", "logo-openai"),
          ("Google", "04/2025", "logo-gemini")]
    tx0 = 0.80
    for i, (nm, dt, lg) in enumerate(tl):
        ex = tx0 + i * 2.0
        add_image(s, ICONS / f"{lg}.png", ex, my + 1.82, 0.34, 0.34)
        text_box(s, ex + 0.42, my + 1.78, 1.5, 0.26, nm, size=11.5, bold=True, color=DEEP)
        text_box(s, ex + 0.42, my + 2.02, 1.5, 0.24, dt, size=11.5, bold=True, italic=True, color=LIGHT)
        if i < 2:
            text_box(s, ex + 1.72, my + 1.82, 0.22, 0.3, "→", size=13, bold=True, color=LIGHT)
    # trust warning — P1 fix (issue #157 review): 3 bullets -> 2 most load-bearing
    # (root cause: code/access; concrete attack vector: prompt injection carrier).
    # Retention-policy point dropped here — it's covered on s25's ZDR block.
    ocean_box(s, 7.25, my, 5.55, 2.55, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, 7.52, my + 0.18, 5.1, 0.72,
             "Standardizing the connection ≠ security of what you connect — and it makes the trust problem worse.",
             size=13.5, bold=True, color=TEAL, line_spacing=1.14)
    for i, w in enumerate([
        "code in your environment / access to data",
        "the description enters the context — a carrier for prompt injection"]):
        circle(s, 7.54, my + 1.06 + i * 0.52 + 0.05, 0.11, TEAL)
        text_box(s, 7.76, my + 1.06 + i * 0.52, 5.0, 0.48, w,
                 size=12.5, color=DEEP, line_spacing=1.10)
    gold_callout(s, 0.55, 6.06, 12.25, 0.80,
                 "No mechanism makes the model more reliable — the ladder rule still holds. Ease of connection is not an argument for connecting.",
                 size=13.5)
    footer(s, "Current savings figures and the scale of the MCP ecosystem — see sources.")
    speaker_notes(s, load_notes("s19"))


def build_s21(p):
    """schema_cycle — agent loop plan→act→check→iterate."""
    s = blank(p)
    slide_title(s, "The agent loop: plan → act → check → iterate.", size=26)
    text_box(s, 0.55, 1.14, 12.25, 0.4,
             "An agent is an architecture where the model does not make one pass, but works in a loop, determining the sequence of steps itself.",
             size=13.5, italic=True, color=MID)
    # 4 step cards in a row with arrows + return arrow below
    steps = [
        ("plan", "Plan", "formulates the next step",
         "myopic / looping plan (does not see the accumulated cost)", MID, False),
        ("act", "Act", "calls a tool (function calling)",
         "the tool fails / stalls, and there is no branch for it", MID, False),
        ("check", "Check", "is the goal reached, is the result correct",
         "validation against an EXTERNAL criterion — not the model's self-assessment (echo of the CoT ceiling)", GOLD, True),
        ("iter", "Iterate", "the loop repeats",
         "no external limit on iterations / cost / time → runaway loop", MID, False),
    ]
    sy, sh = 1.72, 3.55
    bw = 2.92
    gap = 0.20
    x0 = 0.55
    x = x0
    centers = []
    for i, (k, t, sub, fail, col, isgold) in enumerate(steps):
        if isgold:
            ocean_box(s, x, sy, bw, sh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, sy, bw, sh)
        circle(s, x + bw / 2 - 0.26, sy + 0.24, 0.52,
               (GOLD if isgold else MID))
        text_box(s, x + bw / 2 - 0.26, sy + 0.24, 0.52, 0.52, str(i + 1),
                 size=20, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # START badge on Plan (explicit entry point)
        if i == 0:
            chip(s, x + 0.30, sy + 0.30, 0.96, 0.34, "START",
                 fill=TEAL, color=WHITE, size=10.5)
        text_box(s, x + 0.12, sy + 0.90, bw - 0.24, 0.42, t,
                 size=19, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        text_box(s, x + 0.18, sy + 1.36, bw - 0.36, 0.70, sub,
                 size=12, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.12)
        filled_rect(s, x + 0.18, sy + 2.18, bw - 0.36, 1.22,
                    (GOLD_TINT if isgold else SURFACE),
                    stroke=(GOLD if isgold else SOFT_GREY), stroke_pt=1.0,
                    radius=True, radius_adj=0.10)
        text_box(s, x + 0.28, sy + 2.26, bw - 0.56, 1.06,
                 ("failure mode: " + fail),
                 size=10.5, italic=True, color=(DEEP if isgold else SLATE),
                 align=PP_ALIGN.CENTER, line_spacing=1.12)
        centers.append(x + bw / 2)
        if i < 3:
            right_arrow(s, x + bw + 0.01, sy + 0.62, gap - 0.04, 0.46, fill=LIGHT)
        x += bw + gap
    # Explicit return path: down from Iterate → left across → up into Plan
    box_bottom = sy + sh           # 5.27
    ry = box_bottom + 0.55         # horizontal return rail at 5.82
    last_cx = centers[3]
    first_cx = centers[0]
    # down stub from Iterate
    connector(s, last_cx, box_bottom, last_cx, ry, LIGHT, 2.5)
    # horizontal rail right→left
    connector(s, last_cx, ry, first_cx, ry, LIGHT, 2.5)
    # up stub + arrowhead back into Plan
    connector(s, first_cx, ry, first_cx, box_bottom + 0.20, LIGHT, 2.5)
    up_arrow(s, first_cx - 0.13, box_bottom + 0.02, 0.26, 0.24, fill=LIGHT)
    # loop label centered ON the rail, white-backed chip for legibility
    lbl_w = 4.6
    chip(s, (first_cx + last_cx) / 2 - lbl_w / 2, ry - 0.20, lbl_w, 0.40,
         "⟲  the loop REPEATS — back to the \"Plan\" step", fill=LIGHT, color=WHITE,
         size=12.5)
    footer(s, "Loop patterns (ReAct, Reflexion, Plan-and-Execute) — in the chapter. Designing an agent = designing a defense at each of the 4 steps.")
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """schema_matrix / comparison — Workflow vs Agent (2 columns)."""
    s = blank(p)
    slide_title(s, "Workflow vs agent.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "Predictable task → workflow; unpredictable AND the value justifies a multiple-fold increase → agent.",
             size=14.5, italic=True, color=MID)
    cy, chh = 1.78, 2.75
    cw = 6.05
    # workflow
    ocean_box(s, 0.55, cy, cw, chh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "git-fork", 0.83, cy + 0.22, 0.50, "teal")
    text_box(s, 1.45, cy + 0.24, cw - 0.9, 0.5, "Workflow",
             size=17, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.83, cy + 0.85, cw - 0.56, 0.45,
             "LLM and tools along paths predefined in code",
             size=13, bold=True, color=DEEP, line_spacing=1.12)
    for i, t in enumerate(["the sequence of steps is known in advance",
                           "predictable, auditable",
                           "most reliable production systems are a workflow"]):
        circle(s, 0.83, cy + 1.40 + i * 0.43 + 0.06, 0.11, TEAL)
        text_box(s, 1.06, cy + 1.40 + i * 0.43, cw - 1.3, 0.42, t,
                 size=12.5, color=DEEP, line_spacing=1.05)
    # agent
    rx = 6.75
    ocean_box(s, rx, cy, cw, chh)
    icon(s, "bot", rx + 0.28, cy + 0.22, 0.50, "mid")
    text_box(s, rx + 0.90, cy + 0.24, cw - 0.9, 0.5, "Agent",
             size=17, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.28, cy + 0.85, cw - 0.56, 0.45,
             "the LLM dynamically determines its own process",
             size=13, bold=True, color=DEEP, line_spacing=1.12)
    for i, t in enumerate(["the sequence is not fixed in advance",
                           "many times more tokens than a chat",
                           "lower auditability, higher risk of loops"]):
        circle(s, rx + 0.28, cy + 1.40 + i * 0.43 + 0.06, 0.11, MID)
        text_box(s, rx + 0.51, cy + 1.40 + i * 0.43, cw - 1.3, 0.42, t,
                 size=12.5, color=DEEP, line_spacing=1.05)
    # diagnostic question
    ocean_box(s, 0.55, 4.70, 12.25, 1.05, fill=SURFACE, stroke=LIGHT)
    text_box(s, 0.85, 4.80, 11.65, 0.9,
             "Can I write out the sequence of steps in advance, before launch?  yes (even with branches) → workflow  ·  fundamentally no AND the value justifies the multiple-fold cost/risk → agent",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.18)
    # #231: workflow↔agent nesting — the norm, not a third architecture.
    ocean_box(s, 0.55, 5.90, 8.55, 0.98, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, 0.80, 5.98, 8.05, 0.82,
             "Nesting is the norm: a code-review agent calls a \"lint→test→format\" workflow as one step; a ticket-processing workflow delegates parsing a free-form complaint to a mini-agent. Dynamics — only where it is needed.",
             size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
    ocean_box(s, 9.30, 5.90, 3.50, 0.98, fill=SURFACE, stroke=SOFT_GREY, stroke_pt=1.0)
    text_box(s, 9.50, 5.98, 3.10, 0.82,
             "Find the simplest option. \"Too lazy to formalize\" does not make a task unpredictable.",
             size=11.5, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.10)
    speaker_notes(s, load_notes("s22"))


def build_s22b(p):
    """NEW (§4.4) — "what an assistant agent is made of": a map of 5 equipment
    slots. ASSERTION in the title: each slot is a trade-off, not a default
    upgrade (rhymes with the ladder)."""
    s = blank(p)
    slide_title(s, "Every agent equipment slot is a trade-off, not a win.", size=23)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "A real assistant agent (Claude Code, Cursor, Aider) is a plan→act→check→iterate loop PLUS equipment. Five typical slots:",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    slots = [
        ("brain-circuit", "Memory", "what it remembers between sessions"),
        ("file-text", "Instruction rules", "convention files + a task journal"),
        ("puzzle", "Skills", "reusable procedures"),
        ("users", "Subagents", "delegation + isolation"),
        ("cable", "Access / MCP", "access to external tools"),
    ]
    cw = 2.40
    gap = 0.10
    x0 = 0.55
    y = 1.78
    chh = 2.55
    for i, (ic, t, sub) in enumerate(slots):
        x = x0 + i * (cw + gap)
        ocean_box(s, x, y, cw, chh)
        filled_rect(s, x + cw / 2 - 0.42, y + 0.26, 0.84, 0.84, MID, radius=True, radius_adj=0.18)
        icon(s, ic, x + cw / 2 - 0.32, y + 0.36, 0.64, "white")
        text_box(s, x + 0.10, y + 1.28, cw - 0.20, 0.5, t,
                 size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(s, x + 0.14, y + 1.80, cw - 0.28, 0.65, sub,
                 size=11.5, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.12)
        # slot number badge
        circle(s, x + 0.14, y + 0.14, 0.32, GOLD)
        text_box(s, x + 0.14, y + 0.14, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 4.66, 12.25, 1.30,
                 "Just as you don't climb the architecture ladder without a task requirement — don't add memory, subagents or MCP to an agent \"just in case\". Each slot carries its own cost: operational complexity, a new failure surface, a new trust boundary — and it must answer a concrete trigger (Section 5).",
                 size=15)
    speaker_notes(s, load_notes("s22b"))


def build_s22c(p):
    """NEW (§4.5) — agent memory: flat file → mem0/Cognee/Graphiti-Zep.
    Explicit callback to RAG in Section 2 — the same knowledge-scale question."""
    s = blank(p)
    slide_title(s, "Agent memory — the same scale question as RAG.", size=24)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Memory — what the agent remembers BETWEEN sessions (unlike the context of a single conversation). The spectrum runs from a flat file to graph knowledge bases.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # spectrum: flat file -> vector/graph
    sy, sh = 1.86, 2.35
    ocean_box(s, 0.55, sy, 5.55, sh)
    icon(s, "file-text", 0.82, sy + 0.22, 0.50, "mid")
    text_box(s, 1.44, sy + 0.24, 4.4, 0.42, "Flat file", size=16, bold=True,
             color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.82, sy + 0.86, 5.0, 1.35,
             "the agent appends facts to a text log, reads it whole on startup. Works while the log is small and stable — a direct parallel to the RAG criterion \"the corpus fits in the window\".",
             size=13, color=DEEP, line_spacing=1.20)
    right_arrow(s, 6.20, sy + sh / 2 - 0.26, 0.55, 0.52, fill=LIGHT)
    ocean_box(s, 6.95, sy, 5.85, sh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "brain-circuit", 7.22, sy + 0.22, 0.50, "teal")
    text_box(s, 7.84, sy + 0.24, 4.7, 0.42, "Vector / graph knowledge base", size=16,
             bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    for i, m in enumerate([
        "mem0 — user memory between sessions",
        "Cognee — memory on a knowledge graph",
        "Graphiti / Zep — a temporal graph: when a fact is true, when it's stale",
    ]):
        circle(s, 7.24, sy + 0.86 + i * 0.44 + 0.05, 0.11, TEAL)
        text_box(s, 7.46, sy + 0.86 + i * 0.44, 5.1, 0.42, m,
                 size=12.5, color=DEEP, line_spacing=1.08)
    gold_callout(s, 0.55, 4.42, 9.05, 1.90,
                 "The same knowledge-scale question that RAG solves for a document corpus arises here for the agent's own memory. Don't give a graph knowledge base to an agent with short, unrelated sessions — that's the same technical debt without a requirement as RAG for ten articles.",
                 size=14.5)
    # #185: real internet meme "Tuxedo Winnie the Pooh" — a flat file log
    # vs a graph knowledge base for memory (the same spectrum as the diagram top-left).
    pmw, pmh = 3.00, 2.19
    pmx, pmy = 9.80, 4.42
    filled_rect(s, pmx - 0.05, pmy - 0.05, pmw + 0.10, pmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.04)
    add_image(s, WEB / "s22c-pooh-en.png", pmx, pmy, pmw, pmh)
    footer(s, "Source: the public agent-harness-registry (workain lab, independent evaluation).")
    speaker_notes(s, load_notes("s22c"))


def build_s22d(p):
    """NEW (§4.6) — a memory failure (case): Letta Tier D + Anthropic Memory
    Tool Tier B 17%. Letta freshness caveat on the visible layer."""
    s = blank(p)
    slide_title(s, "\"An agent that remembers\" — not always for the better.", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Having memory intuitively seems like a pure win. Independent evaluation shows: sometimes — dramatically not.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # left — Letta case with numbers
    lx, ly, lw, lh = 0.55, 1.86, 6.25, 4.05
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "triangle-alert", lx + 0.26, ly + 0.22, 0.46, "mid")
    text_box(s, lx + 0.84, ly + 0.22, lw - 1.1, 0.42, "Letta — Tier D",
             size=16, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 0.78, lw - 0.52, 0.62,
             "loses to BOTH the bare model AND the flat file on every task. persistbench_v1:",
             size=13, color=DEEP, line_spacing=1.16)
    # mini table of 3 numbers
    rows = [("Bare model", "1.000", "94 s", TEAL),
            ("Flat file", "0.833", "159 s", MID),
            ("Letta", "0.750", "496 s", GOLD)]
    ry = ly + 1.48
    for nm, sc, tm, col in rows:
        filled_rect(s, lx + 0.26, ry, 0.14, 0.34, col)
        text_box(s, lx + 0.52, ry, 2.7, 0.34, nm, size=12.5, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, lx + 3.25, ry, 1.3, 0.34, sc, size=13, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        text_box(s, lx + 4.65, ry, 1.3, 0.34, tm, size=12, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        ry += 0.42
    text_box(s, lx + 0.26, ry + 0.02, lw - 0.52, 0.72,
             "Mechanisms: capitulation under pressure · verbosity drowns the fact · the fact is noticed but not committed.",
             size=12, color=DEEP, line_spacing=1.14)
    # freshness caveat VISIBLE
    filled_rect(s, lx + 0.26, ly + lh - 0.62, lw - 0.52, 0.48, SOFT_GREY,
                radius=True, radius_adj=0.12)
    text_box(s, lx + 0.42, ly + lh - 0.58, lw - 0.80, 0.40,
             "Caveat: the test is on Letta v0.6.7 — behind the current v0.16.8 (~18 mo).",
             size=11, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    # right — Anthropic Memory Tool 17%
    rx, rw = 7.05, 5.75
    ocean_box(s, rx, ly, rw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.28, ly + 0.22, rw - 0.56, 0.42, "Anthropic Memory Tool — Tier B",
             size=16, bold=True, color=DEEP)
    text_box(s, rx + 0.28, ly + 0.70, rw - 0.56, 0.5,
             "A strong result overall — but even it loses data on",
             size=13, color=DEEP, line_spacing=1.14)
    text_runs(s, rx + 0.28, ly + 1.18, rw - 0.56, 0.7, [
        {"text": "17% ", "size": 34, "bold": True, "color": GOLD},
        {"text": "of tasks", "size": 16, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    for i, m in enumerate([
        "an explicit refusal to record an \"ephemeral\" fact",
        "an unmotivated \"out of scope\" refusal",
        "silent lossy compression — the detail can't be recovered",
        "non-reproducibility: the same conversation twice → a different result",
    ]):
        circle(s, rx + 0.30, ly + 2.04 + i * 0.42 + 0.05, 0.11, MID)
        text_box(s, rx + 0.52, ly + 2.04 + i * 0.42, rw - 0.80, 0.42, m,
                 size=12, color=DEEP, line_spacing=1.06)
    footer(s, "agent-harness-registry (workain lab, independent evaluation 2026-07-05). \"Works well\" ≠ \"works always\" — the same lesson as RAG at scale and forgetting.")
    speaker_notes(s, load_notes("s22d"))


def build_s22e(p):
    """NEW (§4.7) — the operational layer: instruction files + presence paradox +
    Honest Lying + claude-code#51735. Intuition diverges from measurement."""
    s = blank(p)
    slide_title(s, "An instruction file for an agent is not a \"magic vaccine\".", size=24)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "The operational layer — convention files (the CLAUDE.md / AGENTS.md class) + a task journal. The intuition \"wrote an instruction → and it got more correct\" diverges from measurement.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # 3 evidence cards
    cards = [
        ("notebook-pen", "presence paradox", "RCT (Gloaguen et al. 2026): the mere presence of an instruction file gives NO significant gain in success — while cost and step count grow.", "Helps only where it actually fills a documentation gap."),
        ("brain-circuit", "Honest Lying", "Dixit, Kamal, Oates 2026: self-authored memory can ENTRENCH a wrong belief — the journal cements an early mistake instead of revising it.", "The reflection is wrong → retries lean on it."),
        ("triangle-alert", "claude-code#51735", "A real case: a past mistake acknowledged in writing did NOT prevent its recurrence 25 days later.", "A record of a failure ≠ a guarantee the behavior changes."),
    ]
    cw, chh = 4.00, 3.35
    cy = 1.80
    x = 0.55
    for i, (ic, t, body, foot) in enumerate(cards):
        isgold = (i == 0)
        if isgold:
            ocean_box(s, x, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, cy, cw, chh)
        icon(s, ic, x + 0.24, cy + 0.24, 0.50, "gold" if isgold else "mid")
        text_box(s, x + 0.86, cy + 0.26, cw - 1.05, 0.5, t,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x + 0.24, cy + 0.92, cw - 0.48, 1.75, body,
                 size=12.5, color=DEEP, line_spacing=1.18)
        filled_rect(s, x + 0.22, cy + chh - 0.86, cw - 0.44, 0.72, SURFACE,
                    stroke=SOFT_GREY, stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, x + 0.36, cy + chh - 0.80, cw - 0.68, 0.62, foot,
                 size=11.5, italic=True, color=(DEEP if isgold else SLATE),
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        x += cw + 0.13
    gold_callout(s, 0.55, 5.30, 9.05, 1.42,
                 "Useful when it fills a real gap; useless when it duplicates what the model already has; can harm when self-authored memory entrenches a mistake. The same pattern as a role in a prompt: \"add X → it gets more correct\" is not confirmed by measurement.",
                 size=13.5)
    # #185: real internet meme "This is fine" — "the instruction file will fix
    # everything" while the agent system is literally on fire (presence paradox).
    fmw, fmh = 3.00, 1.46
    fmx, fmy = 9.80, 5.30
    filled_rect(s, fmx - 0.05, fmy - 0.05, fmw + 0.10, fmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s22e-thisisfine-en.png", fmx, fmy, fmw, fmh)
    speaker_notes(s, load_notes("s22e"))


def build_s23(p):
    """case_study — 3 agent failure cards; compounding chart c23 inside card 2."""
    s = blank(p)
    slide_title(s, "Agent failures.", size=27)
    text_box(s, 0.55, 1.14, 12.25, 0.4,
             "Each failure is a failure mode of the agent loop in a dated case: lesson + alternative.",
             size=14, italic=True, color=MID)
    # card 1 — loop
    c1x, cy, c1w, chh = 0.55, 1.66, 4.05, 4.55
    ocean_box(s, c1x, cy, c1w, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "flame", c1x + 0.22, cy + 0.22, 0.50, "gold")
    text_box(s, c1x + 0.82, cy + 0.24, c1w - 1.0, 0.5, "1. Loop without limits",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, c1x + 0.24, cy + 0.85, c1w - 0.48, 0.95,
             "An agent told to \"sync the orders\" got HTTP 429 → plan→call→429→…",
             size=12, color=DEEP, line_spacing=1.14)
    text_box(s, c1x + 0.24, cy + 1.72, c1w - 0.48, 0.50, "$4,200 in 63 hours",
             size=22, bold=True, color=GOLD, line_spacing=1.0)
    # #233: explicit baseline — a retry script "practically free"
    filled_rect(s, c1x + 0.24, cy + 2.24, c1w - 0.48, 0.62, WHITE, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.12)
    text_box(s, c1x + 0.38, cy + 2.27, c1w - 0.76, 0.56,
             "Baseline: a retry-with-backoff script would solve the same task in seconds-to-minutes, practically free",
             size=10, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    text_box(s, c1x + 0.24, cy + 2.92, c1w - 0.48, 0.56,
             "$4,200 — the price not of automation, but of the wrong architecture choice for a predictable task",
             size=10, italic=True, color=SLATE, line_spacing=1.08)
    text_box(s, c1x + 0.24, cy + 3.54, c1w - 0.48, 0.90,
             "A more fitting architecture: a retry-with-backoff script, not an agent; budget and iteration limits OUTSIDE the agent",
             size=11, bold=True, color=DEEP, line_spacing=1.10)
    # card 2 — compounding (chart)
    c2x, c2w = 4.80, 4.05
    ocean_box(s, c2x, cy, c2w, chh)
    text_box(s, c2x + 0.24, cy + 0.20, c2w - 0.48, 0.4,
             "2. Compounding errors", size=15, bold=True, color=MID)
    add_image(s, CHARTS / "c23-compounding.png", c2x + 0.16, cy + 0.62,
              c2w - 0.32, 2.55)
    text_box(s, c2x + 0.24, cy + 3.30, c2w - 0.48, 1.05,
             "Takeaway: \"tune one step\" is a weak lever; \"fewer transitions + a check between steps\" is a strong one",
             size=12, bold=True, color=DEEP, line_spacing=1.16)
    # card 3 — multi-agent fragility
    c3x, c3w = 9.05, 3.75
    ocean_box(s, c3x, cy, c3w, chh)
    icon(s, "git-fork", c3x + 0.22, cy + 0.22, 0.46, "mid")
    text_box(s, c3x + 0.76, cy + 0.22, c3w - 0.95, 0.55,
             "3. Multi-agent fragility", size=14.5, bold=True, color=MID,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    text_box(s, c3x + 0.24, cy + 0.95, c3w - 0.48, 1.55,
             "dependent subtasks → parallel subagents make conflicting implicit decisions",
             size=12.5, color=DEEP, line_spacing=1.18)
    text_box(s, c3x + 0.24, cy + 2.70, c3w - 0.48, 1.65,
             "A more fitting one: a single-threaded linear agent; multi-agent — only for broadly-parallel independent work",
             size=12.5, bold=True, color=DEEP, line_spacing=1.18)
    footer(s, "Attacks through tools (prompt injection, exfiltration via GitHub MCP) — the 4th class of failures, covered on the security slide. The $4,200 loop — a single author's post-mortem 2026-04 (illustrative); compounding errors — MindStudio 2025–2026.")
    speaker_notes(s, load_notes("s23"))


def build_s25(p):
    """NEW/MERGED (§4.8) — skills + subagents + tool access +
    INTEGRATED security of equal weight (P1: GOLD security block, not a
    caveat below). GitHub MCP heist + ZDR facts live HERE."""
    s = blank(p)
    slide_title(s, "Skills, subagents, access — and the trust boundary.", size=25)
    # top row — 3 equipment slots
    slots = [
        ("puzzle", "Skill", "a reusable “how-to” procedure for a recurring task"),
        ("users", "Subagent", "a separate context window: keep the main one clean + isolate the untrusted"),
        ("cable", "MCP access", "every connection is a new trust boundary and retention policy"),
    ]
    cw, chh = 4.00, 1.55
    cy = 1.12
    x = 0.55
    for ic, t, body in slots:
        ocean_box(s, x, cy, cw, chh)
        icon(s, ic, x + 0.22, cy + 0.22, 0.46, "mid")
        text_box(s, x + 0.80, cy + 0.22, cw - 1.0, 0.42, t,
                 size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.24, cy + 0.74, cw - 0.48, 0.74, body,
                 size=12, color=DEEP, line_spacing=1.14)
        x += cw + 0.13
    # GOLD security block — EQUAL VISUAL WEIGHT (P1), not a footer caveat
    gy, gh = 2.86, 3.60
    ocean_box(s, 0.55, gy, 12.25, gh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=3.0)
    icon(s, "shield-alert", 0.82, gy + 0.22, 0.56, "gold")
    text_box(s, 1.52, gy + 0.24, 11.0, 0.44,
             "Security: the moment an agent delegates and connects, a surface appears that a single call never had",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    # left half — GitHub heist mechanism + 2 ZDR facts
    text_box(s, 0.82, gy + 0.80, 5.85, 0.34, "GitHub MCP: leak via injection (May 2025)",
             size=13.5, bold=True, color=DEEP)
    text_box(s, 0.82, gy + 1.14, 5.85, 0.78,
             "an issue with an embedded instruction + an over-broad token (a PAT for all repos) → the assistant dumped private repositories into a public PR.",
             size=12, color=DEEP, line_spacing=1.16)
    filled_rect(s, 0.82, gy + 1.96, 5.85, 0.46, WHITE, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.14)
    text_box(s, 0.98, gy + 1.99, 5.55, 0.40,
             "Disaster = injection × broad rights. Remove either — the attack fails.",
             size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.82, gy + 2.48, 5.85, 0.44,
             "“We have ZDR” ≠ “stored nowhere”: a court order (NYT v. OpenAI) + third-party services/MCP outside ZDR. Regulated data — don't send without ZDR/BAA.",
             size=11, italic=True, color=DEEP, line_spacing=1.10)
    # right half — 4 rules
    rx = 7.05
    text_box(s, rx, gy + 0.80, 5.5, 0.34, "4 design rules",
             size=13.5, bold=True, color=DEEP)
    rules = [
        ("Least-privilege", "the minimum of tokens/rights"),
        ("Isolate the untrusted", "apart from privileges (a subagent)"),
        ("Human in the loop on writes", "the irreversible — only via a human"),
        ("Allowlist / version pinning", "audited versions; deny by default"),
    ]
    ry = gy + 1.18
    for i, (t, b) in enumerate(rules):
        circle(s, rx, ry + 0.03, 0.30, GOLD)
        text_box(s, rx, ry + 0.03, 0.30, 0.30, str(i + 1),
                 size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 0.42, ry - 0.02, 5.5, 0.30, t, size=13, bold=True, color=DEEP)
        text_box(s, rx + 0.42, ry + 0.26, 5.5, 0.28, b, size=11, color=SLATE)
        ry += 0.58
    speaker_notes(s, load_notes("s25"))


def build_s25b(p):
    """NEW (§4.9) — a survey of real coding agents through the §4.4 equipment
    frame. EXACTLY 4 tools: Claude Code / Aider / Cursor / OpenHands. OpenHands
    marked as an unconfirmed “OpenClaw” hypothesis."""
    s = blank(p)
    slide_title(s, "Real coding agents — through the equipment frame.", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "The difference between tools is not “model quality”, but which equipment slots are filled and where the agent physically lives.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    tools = [
        ("Claude Code", "terminal + IDE · proprietary",
         "broad equipment: memory, instruction files, skills, full subagents, MCP — almost all 5 slots",
         "the price — a lot of operational complexity", MID, False),
        ("Aider", "terminal-first · open source",
         "minimal simplicity: no elaborate memory, no subagents, no skills. ~47k★ on GitHub",
         "a “thin” rig is a deliberate choice, not an underdevelopment", TEAL, False),
        ("Cursor", "desktop IDE (a VS Code fork) · proprietary",
         "the agent lives inside the editor — not a terminal. The form of integration is a separate axis from the rig",
         "where the agent lives is also an architectural decision", MID, False),
        ("OpenHands", "platform on your own server · MIT · ~80k★",
         "broad autonomy, local/Docker/cloud deployment",
         "a working hypothesis by profile match — a likely candidate for “OpenClaw”, not a confirmed fact", GOLD, True),
    ]
    cw, chh = 3.02, 4.05
    cy = 1.74
    x = 0.55
    for nm, meta, body, note, col, ishyp in tools:
        if ishyp:
            ocean_box(s, x, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, cy, cw, chh)
        icon(s, "code", x + 0.22, cy + 0.22, 0.44, "gold" if ishyp else "mid")
        text_box(s, x + 0.76, cy + 0.20, cw - 0.92, 0.48, nm,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x + 0.24, cy + 0.76, cw - 0.48, 0.44, meta,
                 size=10.5, italic=True, color=LIGHT, line_spacing=1.08)
        text_box(s, x + 0.24, cy + 1.22, cw - 0.48, 1.30, body,
                 size=11.5, color=DEEP, line_spacing=1.16)
        nb_h = 1.30 if ishyp else 0.86
        filled_rect(s, x + 0.20, cy + chh - nb_h - 0.14, cw - 0.40, nb_h,
                    WHITE if ishyp else SURFACE, stroke=(GOLD if ishyp else SOFT_GREY),
                    stroke_pt=(1.5 if ishyp else 1.0), radius=True, radius_adj=0.10)
        text_box(s, x + 0.32, cy + chh - nb_h - 0.08, cw - 0.64, nb_h - 0.10, note,
                 size=(9.5 if ishyp else 10), italic=True,
                 color=(DEEP if ishyp else SLATE), line_spacing=1.10)
        x += cw + 0.13
    gold_callout(s, 0.55, 5.98, 12.25, 0.86,
                 "Choosing a tool follows the same rule: don't take the most-equipped one by default — look at which rig your particular task needs.",
                 size=13.5)
    speaker_notes(s, load_notes("s25b"))


def build_s26(p):
    """schema_layered — complexity ladder, bottom-aligned. 6 short rungs +
    trigger label in the gap BELOW each rung (the requirement that opens the
    climb to it). Left col = ladder; right col = rule panel."""
    s = blank(p)
    slide_title(s, "The ladder of architectural complexity.", size=27)
    text_box(s, 0.55, 1.14, 8.40, 0.40,
             "Stay on the bottom rung; climb only when the task requires it.",
             size=13.5, italic=True, color=MID)
    # bottom-up: idx0 = step1 bottom (gold). trig = requirement that opens
    # the climb FROM this rung to the next (shown in the gap above it).
    steps = [
        ("1", "Plain code (no AI)", "you need NL / unstructured input / fuzzy matching", GOLD, True),
        ("2", "One LLM call (prompt; +CoT, +examples in the prompt)", "knowledge is large AND changing AND provenance AND private", MID, False),
        ("3", "RAG / context engineering", "the task is multi-step, the sequence is known in advance", LIGHT, False),
        ("4", "Workflow (predefined paths)", "unpredictable AND the value justifies multiplying cost/risk", LIGHT, False),
        ("5", "Agent (plan→act→check→repeat + limits)", "subtasks are widely parallel AND independent AND high-value", LIGHT, False),
        ("6", "Multi-agent", None, LIGHT, False),
    ]
    n = len(steps)
    step_h = 0.52
    trig_h = 0.34
    bottom_edge = 6.62  # bottom of rung 1
    x0 = 0.55
    full_w = 8.05
    for i, (num, label, trig, col, isgold) in enumerate(steps):
        y = bottom_edge - step_h - i * (step_h + trig_h)
        indent = i * 0.16
        w = full_w - indent
        x = x0 + indent
        if isgold:
            ocean_box(s, x, y, w, step_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, y, w, step_h)
        circle(s, x + 0.12, y + (step_h - 0.34) / 2, 0.34, col)
        text_box(s, x + 0.12, y + (step_h - 0.34) / 2, 0.34, 0.34, num,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.58, y, w - 0.72, step_h, label,
                 size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        if trig:
            # trigger shown in the gap ABOVE this rung (opens climb to next)
            tgy = y - trig_h - 0.01
            text_box(s, x + 0.30, tgy, full_w - 0.50, trig_h,
                     "↑  " + trig, size=10, italic=True, color=SLATE,
                     anchor=MSO_ANCHOR.MIDDLE)
    # PA-1 (owner-approved): direction-of-scale strip alongside the ladder —
    # “more complex ↑” at top / “simpler ↓” at bottom (kills “higher=better” mis-read).
    up_arrow(s, 8.66, 1.98, 0.26, 4.62, fill=COVER_OUTLINE)
    text_box(s, 8.30, 1.58, 1.00, 0.32, "harder ↑", size=10.5, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    text_box(s, 8.30, 6.62, 1.00, 0.32, "simpler ↓", size=10.5, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    # rule panel right (full height, gold — the central rule)
    ocean_box(s, 9.15, 1.55, 3.65, 5.05, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "milestone", 9.42, 1.85, 0.62, "gold")
    text_box(s, 9.42, 2.70, 3.15, 1.75,
             "Stay on the lowest rung that meets the requirements.",
             size=16, bold=True, color=DEEP, line_spacing=1.24)
    text_box(s, 9.42, 4.45, 3.15, 2.05,
             "Every climb is a TRADE-OFF (capability ↔ cost, latency, auditability, attack surface), not a pure win.",
             size=13.5, color=DEEP, line_spacing=1.22)
    footer(s, "The bottom rung is “plain code, no AI”: the ladder begins with the question “do we even need AI at all”.")
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """NEW (§5.2) — “The decision plan”: an 8-step route of questions top to
    bottom (a flowchart with yes/no branches), REPLACING the old 7×7 matrix.
    The bottom priority plate (deterministic → code, STOP) — gold weight."""
    s = blank(p)
    slide_title(s, "The decision plan: a route of questions, not a sum of points.", size=24)
    text_box(s, 0.55, 1.02, 8.4, 0.36,
             "Walk the task top to bottom; stop at the first question that fires. “Yes” → outcome on the right; “no” → next question below.",
             size=12, italic=True, color=MID, line_spacing=1.08)
    # left — vertical routed flow (question -> yes-outcome)
    fx = 0.55
    qw = 5.35           # question box width
    ow = 3.05           # outcome box width
    ax = fx + qw + 0.12  # outcome x
    steps = [
        ("Deterministic and verifiable?", "yes → plain code · STOP", True),
        ("Does one call close it (+CoT)?", "yes → prompt · STOP", False),
        ("Need source provenance?", "yes → RAG grounded in a source / code", False),
        ("Knowledge changes / provenance?", "yes → RAG (not FT)", False),
        ("Need behavior / tone / format?", "yes → fine-tuning (PEFT)", False),
        ("Multi-step, order known?", "yes → workflow · no → agent+limits", False),
        ("Subtasks parallel+independent?", "yes → multi-agent · no → linear", False),
    ]
    y = 1.42
    qh = 0.50
    vg = 0.10
    for i, (q, out, isgold) in enumerate(steps):
        if isgold:
            ocean_box(s, fx, y, qw, qh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.5)
        else:
            ocean_box(s, fx, y, qw, qh)
        circle(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, GOLD if isgold else MID)
        text_box(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, fx + 0.54, y, qw - 0.68, qh, q,
                 size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        # outcome to the right
        filled_rect(s, ax, y + 0.05, ow, qh - 0.10, GOLD if isgold else TEAL_TINT,
                    stroke=(GOLD if isgold else TEAL), stroke_pt=1.2, radius=True, radius_adj=0.12)
        text_box(s, ax + 0.12, y + 0.05, ow - 0.24, qh - 0.10, out,
                 size=10.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        # connector question->outcome (yes →)
        right_arrow(s, fx + qw + 0.005, y + qh / 2 - 0.09, 0.11, 0.18, fill=LIGHT)
        # thin down-connector to next question (no ↓), centred in the gap
        if i < len(steps) - 1:
            connector(s, fx + 0.28, y + qh, fx + 0.28, y + qh + vg, LIGHT, 1.6)
        y += qh + vg
    # step 8 — same row pattern as 1–7 (unified), spanning question+outcome cols
    filled_rect(s, fx, y, qw + 0.12 + ow, qh, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.12)
    circle(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, GOLD)
    text_box(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, 0.32, "8",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, fx + 0.54, y, qw + 0.12 + ow - 0.68, qh,
             "Data sensitive? → at EVERY step: data map + least-privilege + ZDR/BAA",
             size=10.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    # right — worked example + mini-apply
    rx, rw = 9.10, 3.70
    ocean_box(s, rx, 1.46, rw, 2.55, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.24, 1.60, rw - 0.48, 0.34, "Example (task A)", size=13.5,
             bold=True, color=TEAL)
    text_box(s, rx + 0.24, 1.96, rw - 0.48, 1.95,
             "“2000 regulations, changing weekly, an answer citing the clause” → question 4 (changing+provenance) is decisive → RAG with strict grounding in a source. Not fine-tuning (goes stale), not code (NL is needed).",
             size=12, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, 4.12, rw, 1.85)
    text_box(s, rx + 0.24, 4.26, rw - 0.48, 0.34, "Warm-up (task B)", size=13.5,
             bold=True, color=MID)
    text_box(s, rx + 0.24, 4.62, rw - 0.48, 1.30,
             "“a bot over ~150 FAQ, changing once a quarter” — walk the route yourself; we go over it in the seminar.",
             size=12, color=DEEP, line_spacing=1.18)
    # bottom priority plate — gold, the most important line
    gold_callout(s, 0.55, 6.06, 12.25, 0.82,
                 "Route priority: if the task is deterministic and verifiable — plain code, STOP here. AI would only add nondeterminism, cost, latency, and a surface for prompt injection.",
                 size=14)
    speaker_notes(s, load_notes("s27"))


def build_s27b(p):
    """NEW (§5.2b) — “The agent starter kit and when to complicate it”.
    A growth-ladder playbook. Rhymes with the s26 ladder and the 5-slot map §4.4."""
    s = blank(p)
    slide_title(s, "The agent starter kit — and when to complicate it.", size=24)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "The same “don't complicate without a requirement” ladder, applied one level down — to a single agent's rig, not the whole system.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # left — thin default
    lx, ly, lw, lh = 0.55, 1.86, 4.45, 4.05
    ocean_box(s, lx, ly, lw, lh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "package", lx + 0.26, ly + 0.24, 0.54, "teal")
    text_box(s, lx + 0.92, ly + 0.26, lw - 1.1, 0.5, "By default — a thin agent",
             size=16, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    for i, t in enumerate([
        "one instruction file",
        "flat memory",
        "NO subagents",
        "a minimal set of skills",
        "minimal MCP access",
    ]):
        circle(s, lx + 0.30, ly + 1.02 + i * 0.48 + 0.06, 0.12, TEAL)
        text_box(s, lx + 0.56, ly + 1.02 + i * 0.48, lw - 0.84, 0.44, t,
                 size=14, color=DEEP, line_spacing=1.1)
    text_box(s, lx + 0.28, ly + lh - 0.62, lw - 0.56, 0.56,
             "The burden of proof is on complication, not on simplicity.",
             size=12, italic=True, color=DEEP, line_spacing=1.14)
    # right — 3 justified triggers (narrowed — the far-right column is under the meme)
    rx, rw = 5.30, 4.30
    triggers = [
        ("brain-circuit", "A memory backend — when:",
         "history has outgrown the context OR you need to search over facts. The same criterion as prompt→RAG."),
        ("users", "Subagents — when:",
         "a subtask needs a separate window OR isolation of untrusted work (least-privilege)."),
        ("cable", "More MCP access — when:",
         "a specific task needs a specific tool — not “just in case”. Every connection is a trust boundary."),
    ]
    ty2 = 1.86
    th = 1.24
    for ic, t, b in triggers:
        ocean_box(s, rx, ty2, rw, th)
        icon(s, ic, rx + 0.22, ty2 + 0.20, 0.42, "mid")
        text_box(s, rx + 0.76, ty2 + 0.16, rw - 0.94, 0.36, t,
                 size=13.5, bold=True, color=MID)
        text_box(s, rx + 0.76, ty2 + 0.50, rw - 0.94, 0.68, b,
                 size=11.5, color=DEEP, line_spacing=1.12)
        ty2 += th + 0.16
    # #185: the real “Two guys on a bus” internet meme — sad (complicated it
    # just in case) vs happy (started with a thin agent). Far right.
    bmx, bmy, bmw, bmh = 9.75, 2.35, 3.05, 3.10
    ocean_box(s, bmx, bmy, bmw, bmh)
    add_image(s, WEB / "s27b-bus-en.png", bmx + 0.12, bmy + 0.12,
              bmw - 0.24, bmh - 0.24)
    gold_callout(s, 0.55, 6.06, 12.25, 0.82,
                 "The same principle as the architecture ladder — applied to a single agent's rig: the presence paradox showed that even an instruction file “as a ritual” doesn't work. Complicate against a specific verifiable trigger.",
                 size=13)
    speaker_notes(s, load_notes("s27b"))


def build_s29(p):
    """assertion_visual — human validator + MIT NANDA ~95% donut (c29)."""
    s = blank(p)
    slide_title(s, "Human validator + the MIT NANDA lesson.", size=27)
    # left — human validator
    lx, ly, lw, lh = 0.55, 1.40, 6.55, 4.45
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "user-check", lx + 0.26, ly + 0.24, 0.52, "mid")
    text_box(s, lx + 0.92, ly + 0.24, lw - 1.1, 0.50,
             "The agent acts — a human checks the result and the facts",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, lx + 0.28, ly + 0.80, lw - 0.56, 0.56,
             "against an independent source of truth, not by the plausibility of the reasoning: the model's self-explanation ≠ control.",
             size=12, italic=True, color=SLATE, line_spacing=1.14)
    # #237: three dimensions of the human-validator role
    dims = [
        ("Degree of autonomy", "from “a human presses the button” to “the agent notifies after the fact”"),
        ("Trust scope", "reads (easy to roll back) vs writes · reversible vs irreversible"),
        ("Continuous monitoring", "quality metrics all the time, not a one-off check at the start"),
    ]
    py = ly + 1.52
    for t, b in dims:
        filled_rect(s, lx + 0.28, py, lw - 0.56, 0.90, SURFACE, stroke=SOFT_GREY,
                    stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, lx + 0.44, py + 0.10, lw - 0.86, 0.34, t,
                 size=13.5, bold=True, color=MID)
        text_box(s, lx + 0.44, py + 0.44, lw - 0.86, 0.42, b,
                 size=11.5, color=DEEP, line_spacing=1.10)
        py += 0.98
    # right — NANDA donut
    rx, rw = 7.35, 5.45
    ocean_box(s, rx, ly, rw, 3.10, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    add_image(s, CHARTS / "c29-nanda.png", rx + 0.20, ly + 0.20, 2.45, 2.65)
    text_box(s, rx + 2.75, ly + 0.35, rw - 3.0, 0.95, "~95%",
             size=42, bold=True, color=GOLD, line_spacing=1.0)
    text_box(s, rx + 2.75, ly + 1.30, rw - 3.0, 1.65,
             "of enterprise GenAI pilots with no measurable return on investment — the root is in the learning gap and the integration failure, not model quality",
             size=12.5, color=DEEP, line_spacing=1.18)
    ocean_box(s, rx, ly + 3.25, rw, 1.20)
    text_box(s, rx + 0.26, ly + 3.38, rw - 0.52, 0.95,
             "“Launch AI” ≠ “get value”. Architectural-integration discipline decides. Sometimes the right answer is the simplest architecture or no AI.",
             size=12, bold=True, color=DEEP, line_spacing=1.16,
             anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "MIT NANDA, State of AI in Business 2025 — a report with a methodology (150 interviews + 350 survey + 300 deployments), not a universal law.")
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """hero_closing / bridge — the bridge to Lecture 4 “AI in software
    development (the SDLC × AI axis)” (issue #170). Left: 4 stage boxes SDLC × AI,
    each with a teal anchor from L3. Right: a hero dev photo ≥40%. Gold bridge title."""
    s = blank(p)
    hx, hy, hw, hh = 8.10, 0.0, 5.233, 7.5
    hero_image(s, SCREENSHOTS / "s30-coding.jpg", hx, hy, hw, hh)
    # #185/#317: the photo attribution caption is removed from visible (attribution.md — legal).
    lx, lw = 0.55, 7.30
    text_box(s, lx, 0.42, lw, 0.94,
             "Next — industries. Start: software development, the SDLC × AI axis.",
             size=24, bold=True, color=DEEP, line_spacing=1.06)
    text_box(s, lx, 1.34, lw, 0.40,
             "Lecture 4 takes the same apparatus and lays it out across the software lifecycle stages. Each stage leans on an anchor from this lecture:",
             size=12.5, italic=True, color=MID, line_spacing=1.14)
    # 2x2 SDLC × AI boxes, each with a teal anchor line from Lec 3
    boxes = [
        ("Requirements / design", "picking an architecture per step → the 6-rung ladder"),
        ("Coding", "coding agents through the equipment frame; failures as a CLASS"),
        ("Testing", "the trust boundary and least-privilege; data outside ZDR"),
        ("Operations", "a human validator checks the result, not the self-explanation"),
    ]
    bx0, by0 = lx, 1.90
    bw, bh = 3.55, 1.62
    gap = 0.15
    for i, (t, anchor) in enumerate(boxes):
        x = bx0 + (i % 2) * (bw + gap)
        y = by0 + (i // 2) * (bh + gap)
        ocean_box(s, x, y, bw, bh)
        chip(s, x + 0.18, y + 0.16, min(bw - 0.36, 0.15 * len(t) + 0.4), 0.38,
             t, fill=MID, color=WHITE, size=11.5)
        text_box(s, x + 0.22, y + 0.66, bw - 0.44, 0.86, anchor, size=11.5,
                 italic=True, color=TEAL, line_spacing=1.16)
    gold_callout(s, lx, 5.42, lw, 0.72,
                 "Assignment — Seminar 3: run the architecture-choice checklist on 3 cases (chat / agent / RAG / API).",
                 size=13)
    gold_callout(s, lx, 6.28, lw, 0.78,
                 "What to do: this frame is the base for Lectures 4–17. On each industry lecture, run the same choice checklist.",
                 size=13)
    speaker_notes(s, load_notes("s30"))


# ============================================================
# v3 new builders (suffix-ID, plan §4 U-1…U-7) — NO renumber s01–s30.
# ============================================================
def build_s04a(p):
    """section_divider — Section 1 “The prompt and its boundaries” (U-1)."""
    build_section_divider(
        p, 1, "Section 1", "The prompt and its boundaries",
        "We saw the whole ladder — now from the bottom: what one call can do and where its ceiling is, before complicating anything.",
        "s04a",
        image_src=WEB / "div-r1-knife.jpg",
        tag="one precise cut · 3 cases · 1 failure")


def build_s13a(p):
    """section_divider — Section 3 “Fine-tuning vs prompt vs RAG” (U-3)."""
    build_section_divider(
        p, 3, "Section 3", "Fine-tuning vs prompt vs RAG",
        "We solved the knowledge problem through RAG. But what if the problem is not knowledge, but the model's behavior — its tone, format, policy?",
        "s13a",
        image_src=WEB / "div-r3-tuning.jpg",
        tag="behavior tuning · 4 cases · 1 failure")


def build_s13b(p):
    """assertion_visual — the definition of fine-tuning BEFORE the critique (U-2).

    Definition on top → center: a mini pipeline schema
    [pretrained model]＋[dataset] → fine-tuning → [fine-tuned weights] →
    bottom: a contrast plate CONTEXT vs WEIGHTS. Gold anchor — “WEIGHTS”.
    Schema §5.5 Process/Pipeline checklist.
    """
    s = blank(p)
    slide_title(s, "What fine-tuning is.", size=27)
    text_box(s, 0.55, 1.14, 12.25, 0.74,
             "Fine-tuning is continuing to train an already-built model on your data. In L1 — a type of use; here — an architectural choice, one of the rungs of the ladder.",
             size=14, italic=True, color=MID, line_spacing=1.18)
    # mini-schema pipeline in ocean box
    sy, sh = 1.98, 2.48
    ocean_box(s, 0.55, sy, 12.25, sh)
    # 3 nodes + «+» (n1→n2) + arrow with «fine-tuning» label (n2→n3)
    bw, bh = 2.85, 1.46
    by = sy + 0.66                     # node band lowered → label clearance
    n1x = 1.35                         # centered group (box inner 0.55..12.80)
    plus_x = n1x + bw                  # 4.20 — «+» zone 0.58 wide
    n2x = plus_x + 0.58                # 4.78
    arr_x0 = n2x + bw                  # 7.63 — arrow zone 1.50 wide
    n3x = arr_x0 + 1.50                # 9.13
    # node 1 — pretrained model
    filled_rect(s, n1x, by, bw, bh, SURFACE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    icon(s, "cpu", n1x + bw / 2 - 0.23, by + 0.14, 0.42, "mid")
    text_box(s, n1x + 0.10, by + 0.60, bw - 0.20, 0.56,
             "Pretrained\nmodel", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.02)
    text_box(s, n1x + 0.14, by + 1.18, bw - 0.28, 0.26,
             "general weights", size=11.5, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    # «+» between n1 and n2
    text_box(s, plus_x, by + bh / 2 - 0.32, 0.58, 0.64, "+",
             size=30, bold=True, color=MID, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # node 2 — your dataset
    filled_rect(s, n2x, by, bw, bh, SURFACE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    icon(s, "database", n2x + bw / 2 - 0.23, by + 0.14, 0.42, "teal")
    text_box(s, n2x + 0.14, by + 0.60, bw - 0.28, 0.30,
             "Your dataset", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, n2x + 0.10, by + 0.94, bw - 0.20, 0.50,
             "examples of the\nwanted behavior", size=11.5, italic=True,
             color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.02)
    # «fine-tuning» label ABOVE arrow (clear vertical separation)
    text_box(s, arr_x0 - 0.05, by - 0.42, 1.60, 0.32, "fine-tuning",
             size=13, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # arrow node2 → node3
    right_arrow(s, arr_x0 + 0.06, by + bh / 2 - 0.21, 1.38, 0.42, fill=MID)
    # node 3 — fine-tuned weights (gold = the changed thing)
    filled_rect(s, n3x, by, bw, bh, GOLD_TINT, stroke=GOLD, stroke_pt=2.0,
                radius=True, radius_adj=0.10)
    icon(s, "sliders-horizontal", n3x + bw / 2 - 0.23, by + 0.14, 0.42,
         "gold")
    text_box(s, n3x + 0.10, by + 0.60, bw - 0.20, 0.56,
             "Fine-tuned\nWEIGHTS", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.02)
    text_box(s, n3x + 0.14, by + 1.18, bw - 0.28, 0.26,
             "the model is now different", size=11.5, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    # contrast strip — CONTEXT vs WEIGHTS (2 halves)
    cy, ch = 4.55, 1.55
    ocean_box(s, 0.55, cy, 6.05, ch, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, 0.83, cy + 0.18, 5.50, 0.34, "Prompt / RAG → CONTEXT",
             size=15, bold=True, color=TEAL)
    text_box(s, 0.83, cy + 0.58, 5.55, 0.88,
             "They change only the input — the weights are untouched; the effect lives only within the request.",
             size=13, color=DEEP, line_spacing=1.22)
    ocean_box(s, 6.75, cy, 6.05, ch)
    text_box(s, 7.03, cy + 0.18, 5.50, 0.34, "Fine-tuning → THE WEIGHTS",
             size=15, bold=True, color=MID)
    text_box(s, 7.03, cy + 0.58, 5.55, 0.88,
             "The change is baked into the model — it applies always and costs more than what context changes.",
             size=13, color=DEEP, line_spacing=1.22)
    gold_callout(s, 0.55, 6.28, 12.25, 0.78,
                 "Prompt/RAG = “what to show the model”.  Fine-tuning = “change the model itself”. In practice “fine-tune” almost always means LoRA/PEFT, not retraining all the weights (next slide).",
                 size=14)
    speaker_notes(s, load_notes("s13b"))



def build_s25a(p):
    """section_divider — Section 5 “How to choose: the decision framework” (U-5b)."""
    build_section_divider(
        p, 5, "Section 5", "How to choose: the decision framework",
        "We went through every architecture separately — and where each one fails. Now we pull it together into one choice tool.",
        "s25a",
        image_src=WEB / "div-r5-control-panel.jpg",
        tag="the choice tool · 4 cases")


# ============================================================
# v5b classic-base-first (issue #185 WP8) — 51→56 slides.
# One “classical baseline from scratch” slide per content section
# (§1–§5), inserted right after the section divider, before the AI
# part. Shared template: 3 classical-baseline cards + gold “what to keep
# from the classics” + a bridge to the AI part. NO renumber (mnemonic ids
# s-classic-*). Schema §5.5: 3-column tile, mass-balanced, single-line names.
# ============================================================

def build_classic_base(p, sid, *, title, intro, cards, keep_text, bridge):
    """Shared "classical baseline" slide: title + intro line + 3 tile cards
    (icon + bold single-line name + body) + gold "what to keep" callout +
    teal-tint bridge strip. `cards` = list of 3 (icon_name, name, body)."""
    s = blank(p)
    slide_title(s, title, size=26)
    text_box(s, 0.55, 1.14, 12.25, 0.60, intro,
             size=14, italic=True, color=MID, line_spacing=1.14)
    # 3 tile cards — equal mass, full width. Name may take 2 lines;
    # body — up to 5 lines, all inside the card (no overflow under the gold plate).
    n = len(cards)
    gap = 0.24
    x0 = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    cy, chh = 1.86, 2.80
    for i, (ic, name, body) in enumerate(cards):
        x = x0 + i * (cw + gap)
        ocean_box(s, x, cy, cw, chh)
        icon(s, ic, x + 0.28, cy + 0.24, 0.54, "mid")
        text_box(s, x + 0.28, cy + 0.92, cw - 0.56, 0.72, name,
                 size=15.5, bold=True, color=MID, line_spacing=1.04)
        text_box(s, x + 0.28, cy + 1.66, cw - 0.56, chh - 1.80, body,
                 size=12, color=DEEP, line_spacing=1.14)
    # gold "what to keep from the classics"
    ky, kh = 4.80, 1.36
    filled_rect(s, 0.55, ky, 12.25, kh, GOLD_TINT, stroke=GOLD, stroke_pt=1.75,
                radius=True, radius_adj=0.06)
    text_box(s, 0.83, ky + 0.11, 2.4, 0.30, "WHAT TO KEEP",
             size=12, bold=True, color=GOLD)
    text_box(s, 0.83, ky + 0.44, 11.6, kh - 0.54, keep_text,
             size=13.5, bold=True, color=DEEP, line_spacing=1.14)
    # teal-tint bridge strip to the AI part
    by, bh = 6.26, 0.80
    filled_rect(s, 0.55, by, 12.25, bh, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_runs(s, 0.83, by + 0.11, 11.6, bh - 0.20, [
        {"text": "Bridge: ", "size": 12, "bold": True, "color": TEAL},
        {"text": bridge, "size": 12, "color": DEEP},
    ], line_spacing=1.12, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes(sid))
    return s


def build_s_classic_prompt(p):
    """§1.0 — classical baseline of section 1: precise task specification."""
    build_classic_base(
        p, "s-classic-prompt",
        title="How the task was framed before the prompt — and what that changes.",
        intro="Before large models, \"make the system do what you need\" meant not a wish in "
              "natural language, but a precise specification and a deterministic program.",
        cards=[
            ("file-text", "Precise specification / requirements",
             "Pre- and postconditions, invariants, acceptance criteria (Z notation, TLA+, Design by "
             "Contract). The result is deterministic and verifiable."),
            ("git-fork", "Imperative vs declarative",
             "\"How to do it\" (a step-by-step algorithm) vs \"what to get\" (SQL/Prolog describe the "
             "result, the engine decides how)."),
            ("braces", "Interface contract",
             "A precise input/output agreement (type signatures, OpenAPI, Protobuf). One "
             "correct meaning and a way to check it."),
        ],
        keep_text="The discipline of precise framing: a precise prompt is the same requirements spec in natural language. "
                  "For the deterministic and verifiable (arithmetic, schema validation, "
                  "rule-based routing) — use classical code, not a prompt.",
        bridge="a prompt frames the task for a probabilistic system in natural language; hence both its "
               "power (no need to specify the un-specifiable) and its boundary (no single meaning, "
               "no determinism).",
    )


def build_s_classic_rag(p):
    """§2.0 — classical baseline of section 2: classical information retrieval."""
    build_classic_base(
        p, "s-classic-rag",
        title="How we searched text before embeddings.",
        intro="The R in RAG is retrieval, search: a discipline with half a century of history. Exactly how "
              "it works determines where RAG works and where it breaks.",
        cards=[
            ("book-open", "Inverted index",
             "For each word — a list of documents where it occurs (a machine catalog). "
             "The foundation of Lucene, Elasticsearch, PostgreSQL full-text."),
            ("route", "Boolean search",
             "A query as a logical expression (\"error AND authentication NOT tomcat\"): precise, "
             "predictable, explainable selection."),
            ("list-ordered", "TF-IDF → BM25",
             "Ranking by word importance (rarer in the collection — stronger signal). BM25 (Okapi) — "
             "a cheap, explainable baseline many do not beat."),
        ],
        keep_text="The classics are precise on codes and identifiers, where semantic search blurs. A strong "
                  "RAG-2026 is a hybrid of BM25 + dense vectors, lexical filters over metadata, "
                  "ranking discipline (a reranker) and observability: recall/precision on "
                  "a golden set.",
        bridge="semantic search over embeddings (Lecture 2) adds semantic "
               "matching on top of the classics instead of lexical — but does not replace it. RAG extends classical "
               "search, it does not abolish it.",
    )


def build_s_classic_ft(p):
    """§3.0 — classical baseline of section 3: classical machine learning."""
    build_classic_base(
        p, "s-classic-ft",
        title="How the ML task was solved before large models.",
        intro="Fine-tuning a large model is not an exotic novelty of a new era, but a direct continuation "
              "of the classical transfer-learning scheme. Let's rebuild it from scratch.",
        cards=[
            ("database", "Sample + ground truth",
             "A labeled set \"input → correct answer\" (ground truth, reference labeling); "
             "the model learns to reproduce it and generalize to new data."),
            ("scale", "Train / val / test split",
             "Three disjoint parts: you train on train, tune on validation, measure once on test. "
             "The rule: you must not test on the training data."),
            ("git-branch", "Transfer learning",
             "Take a model pretrained on a large corpus and cheaply fine-tune it for your narrow "
             "task: faster and more accurate than from scratch (\"pretraining → fine-tuning\")."),
        ],
        keep_text="Eval sets (golden set) — without them catastrophic forgetting is invisible; "
                  "versioning of data and weights for rollback; train/test discipline against leakage; "
                  "drift monitoring. LoRA made the fine-tuning step cheaper, but not the discipline around it.",
        bridge="PEFT/LoRA is the same transfer learning, pushed to the limit of cheapness and on top of "
               "an incomparably larger model. The idea is not new — what became new is the scale of the base "
               "model and the cost of the step.",
    )


def build_s_classic_agents(p):
    """§4.0 — classical baseline of section 4: classical automation."""
    build_classic_base(
        p, "s-classic-agents",
        title="Controlled automation existed long before agents.",
        intro="\"Agent\" sounds like an invention of the large-model era, but controlled process "
              "automation is a mature discipline. Without it you can't judge what an agent adds and what it breaks.",
        cards=[
            ("waypoints", "Finite-state machine",
             "A set of states and event-driven transition rules (\"created → in progress → closed\"). "
             "Possible transitions are visible, impossible ones are excluded by construction."),
            ("git-branch", "Workflow engines",
             "Execute a predefined process by a fixed schema: BPMN, DAG orchestrators "
             "(Airflow), RPA. The order of steps is defined in advance, not on the fly."),
            ("route", "Control loop",
             "plan → act → check as a feedback loop (control theory, SCADA); the military "
             "analog is the OODA loop. Not a single action, but a loop with correction."),
        ],
        keep_text="Deterministic workflows where possible; idempotency (repeating a step does not break "
                  "state); the least-privilege principle for tools; "
                  "limits and audit of the loop. Reliability practices of the control loop don't go anywhere.",
        bridge="an LLM agent is the same plan→act→check control loop, but the steps in it are generated by the model "
               "non-deterministically, rather than choosing a fixed automaton rule. Hence both the flexibility "
               "and the failures of the section.",
    )


def build_s_classic_framework(p):
    """§5.0 — classical baseline of section 5: classical technology choice."""
    build_classic_base(
        p, "s-classic-framework",
        title="How an engineer chose technology before the AI hype.",
        intro="The architecture choice itself is not a new AI procedure, but a direct application of classical "
              "engineering-decision principles. The section's ladder and checklist stand on them.",
        cards=[
            ("clipboard-list", "From the requirement, not the tool",
             "Requirements engineering: first pin down what the system must do, then "
             "pick the tool. Build-vs-buy: build your own or take a ready one."),
            ("scale", "KISS + YAGNI",
             "Choose the simplest solution that meets the requirement; don't build in capacity \"for "
             "the future\" until a concrete requirement demands it."),
            ("milestone", "Least power",
             "From the W3C notes (Berners-Lee and Mendelsohn): take the least powerful of the sufficient "
             "tools — it's easier to analyze, verify and maintain."),
        ],
        keep_text="The simplest sufficient architecture by default, and the burden of proof is on "
                  "whoever wants to complicate it. In the AI era the principle isn't abolished, it gets costlier: an extra "
                  "step adds non-determinism, token cost, latency, attack surface.",
        bridge="the ladder of architectures (code → single call → RAG → workflow → agent → multi-agent) is the "
               "same least-power principle over AI architectures: stay on the lowest sufficient "
               "step, climb only to meet a requirement.",
    )


def build_s01b(p):
    """case_study — Air Canada as the first walkthrough of the class "wrong
    architecture for the task". Hero photo of the 787 on the left ≥40%, the case chronicle +
    takeaway on the right. (Reuse of the v4 s01 hero photo; the s01 meme hook comes before this.)"""
    s = blank(p)
    hx, hy, hw, hh = 0.0, 0.0, 6.05, 7.5
    hero_image(s, SCREENSHOTS / "s01-aircanada.jpg", hx, hy, hw, hh)
    filled_rect(s, 0.0, 7.10, 3.9, 0.40, DEEP)
    text_box(s, 0.14, 7.12, 3.7, 0.34,
             "Photo: Air Canada Boeing 787 · Wikimedia · CC-BY-SA",
             size=10.5, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    rx = 6.45
    text_box(s, rx, 0.50, 6.35, 1.50,
             "First walkthrough: the chatbot invented a policy — the company pays.",
             size=27, bold=True, color=DEEP, line_spacing=1.06)
    text_box(s, rx, 2.05, 6.35, 0.36,
             "Moffatt v. Air Canada · BC tribunal · 02/14/2024",
             size=13.5, bold=True, color=TEAL)
    chron = [
        "A passenger asked the chatbot about a bereavement fare",
        "Bot: \"buy at full price, get the difference back within 90 days\"",
        "The real policy did not allow this — and it was on the very page the bot linked to",
        "Tribunal: \"the bot is not a separate legal entity\" → the company refunded $812.02",
    ]
    cy = 2.52
    row_h = 0.86
    for i, t in enumerate(chron):
        circle(s, rx, cy + 0.03, 0.32, MID if i < 3 else GOLD)
        text_box(s, rx, cy + 0.03, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=(WHITE if i < 3 else DEEP),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 0.50, cy, 5.85, row_h - 0.04, t,
                 size=13, color=DEEP, line_spacing=1.10,
                 anchor=MSO_ANCHOR.MIDDLE)
        cy += row_h
    gold_callout(s, rx, 5.98, 6.35, 0.90,
                 "What to do: a legally binding answer to a customer is not for a generative bot. You need a deterministic source of policy; the bot is only a navigator to the official document.",
                 size=12.5)
    speaker_notes(s, load_notes("s01b"))


def build_s05c(p):
    """assertion_visual (§1.2) — two meanings of the word "role": role-persona (tone)
    ≠ protocol roles system/user/assistant (dialogue structure). system priority
    is a tendency (~63.8%), not a boundary; STI/ChatInject spoof the role.
    Left — chat template, right — STI injection; number plates; gold takeaway."""
    s = blank(p)
    slide_title(s, "\"Role\" is two different mechanisms. Don't confuse them.", size=26)
    text_box(s, 0.55, 1.16, 12.25, 0.44,
             "A role-persona (\"you are a lawyer\") tunes the tone. The protocol role system/user/assistant is markup of \"who is speaking\" in the token stream. The second is constantly mistaken for the first.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # LEFT — chat template assembly
    lx, ly, lw, lh = 0.55, 1.82, 6.05, 2.94
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "braces", lx + 0.26, ly + 0.20, 0.44, "mid")
    text_box(s, lx + 0.84, ly + 0.20, lw - 1.0, 0.42,
             "Chat template: a list of roles → a flat token stream",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.02)
    # mono token strip
    filled_rect(s, lx + 0.26, ly + 0.86, lw - 0.52, 0.94, DEEP, radius=True,
                radius_adj=0.06)
    text_runs(s, lx + 0.42, ly + 0.94, lw - 0.84, 0.80, [
        {"text": "<|im_start|>system", "size": 12, "color": GOLD,
         "font": FONT_MONO, "bold": True},
        {"text": "  behavior rules  ", "size": 12, "color": WHITE,
         "font": FONT_MONO},
        {"text": "<|im_end|>", "size": 12, "color": TEAL, "font": FONT_MONO},
        {"text": "\n<|im_start|>user", "size": 12, "color": LIGHT,
         "font": FONT_MONO, "bold": True, "newpara": True},
        {"text": "  question  ", "size": 12, "color": WHITE, "font": FONT_MONO},
        {"text": "<|im_end|>", "size": 12, "color": TEAL, "font": FONT_MONO},
    ], line_spacing=1.2)
    text_box(s, lx + 0.26, ly + 1.94, lw - 0.52, 0.90,
             "Special markup tokens (ChatML, Llama 4) are assembled by the template — it can be read and spoofed.",
             size=12, color=DEEP, line_spacing=1.16)
    # RIGHT — priority is a tendency, not a guarantee + STI
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, 1.36, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.40,
             "system priority — a tendency, not a boundary",
             size=14, bold=True, color=DEEP)
    text_runs(s, rx + 0.24, ly + 0.58, rw - 0.48, 0.70, [
        {"text": "GPT-4o obeys the priority ~", "size": 13, "color": DEEP},
        {"text": "63.8%", "size": 16, "bold": True, "color": MID},
        {"text": "  — still ≠ 100%: it's a tendency.", "size": 13, "color": DEEP},
    ], line_spacing=1.16)
    ocean_box(s, rx, ly + 1.48, rw, 1.46, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    icon(s, "shield-alert", rx + 0.24, ly + 1.66, 0.42, "teal")
    text_box(s, rx + 0.80, ly + 1.66, rw - 1.0, 0.40,
             "STI / role spoofing: the role is faked",
             size=14, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_runs(s, rx + 0.24, ly + 2.12, rw - 0.48, 0.74, [
        {"text": "The string ", "size": 12, "color": DEEP},
        {"text": "<|im_start|>assistant", "size": 11.5, "color": MID,
         "font": FONT_MONO, "bold": True},
        {"text": " from external text → ChatInject ASR ", "size": 12,
         "color": DEEP},
        {"text": "5.18%→32.05%", "size": 14, "bold": True, "color": TEAL},
        {"text": " (Llama-4 up to 88.3%).", "size": 12, "color": DEEP},
    ], line_spacing=1.14)
    gold_callout(s, 0.55, 4.92, 9.35, 1.58,
                 "What to do: don't design defenses assuming \"the system always outranks the user\". Escape special tokens in INCOMING external content and check the chat template of local models — a foreign template quietly breaks the priority.",
                 size=13.5)
    # #185: the real internet meme "Is this a pigeon?" — the system protocol
    # role is mistaken for a hard boundary (it's only a tendency, not a guarantee).
    gmw, gmh = 2.66, 1.78
    gmx, gmy = 10.05, 4.92
    filled_rect(s, gmx - 0.05, gmy - 0.05, gmw + 0.10, gmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s05c-pigeon-en.png", gmx, gmy, gmw, gmh)
    footer(s, "The full walkthrough of prompt injection as a class of agent attacks is in the agents section (security).")
    speaker_notes(s, load_notes("s05c"))


def build_s07(p):
    """case_study (§1.5) — the limit of CoT: faithfulness. CoT is generated
    text, not a protocol. A hint experiment; Claude 3.7 ~25% / R1 ~39%;
    worse on hard tasks. Check the result, not the self-explanation."""
    s = blank(p)
    slide_title(s, "Reasoning out loud is not an audit log.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.46,
             "Faithfulness (of the explanation) — how well the spoken chain reflects the real cause of the answer. It was measured directly — and it is low.",
             size=14, italic=True, color=MID, line_spacing=1.14)
    # LEFT — the experiment (2 runs)
    lx, ly, lw = 0.55, 1.86, 6.05
    ocean_box(s, lx, ly, lw, 3.30)
    text_box(s, lx + 0.24, ly + 0.16, lw - 0.48, 0.40,
             "Anthropic experiment (Apr. 2025)", size=14, bold=True,
             color=DEEP)
    runs = [
        ("The task is given twice: without a hint and with a hint in the prompt that changes the answer (\"the professor thinks it's C\")", "cpu"),
        ("They look: when the answer changed under the hint — did the model mention it in its reasoning?", "eye-off"),
        ("Often the model changes the answer, but builds a DIFFERENT, invented rationale, without naming the real cause", "triangle-alert"),
    ]
    ry = ly + 0.68
    for t, ic in runs:
        icon(s, ic, lx + 0.26, ry + 0.04, 0.40, "mid")
        text_box(s, lx + 0.82, ry, lw - 1.06, 0.86, t,
                 size=12.5, color=DEEP, line_spacing=1.14)
        ry += 0.92
    # RIGHT — the numbers
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, 3.30, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.40,
             "How often the model admits the hint", size=14, bold=True,
             color=TEAL)
    bars = [("Claude 3.7 Sonnet", 25, MID), ("DeepSeek R1", 39, TEAL)]
    bt = ly + 0.78
    max_w = rw - 1.9
    for nm, pct, col in bars:
        text_box(s, rx + 0.24, bt, rw - 0.48, 0.30, nm, size=13, bold=True,
                 color=DEEP)
        filled_rect(s, rx + 0.24, bt + 0.34, max_w, 0.40, SOFT_GREY,
                    radius=True, radius_adj=0.3)
        filled_rect(s, rx + 0.24, bt + 0.34, max_w * pct / 100.0, 0.40, col,
                    radius=True, radius_adj=0.3)
        text_box(s, rx + 0.24 + max_w + 0.12, bt + 0.30, 1.4, 0.48,
                 f"~{pct}%", size=18, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        bt += 0.92
    text_box(s, rx + 0.24, bt + 0.02, rw - 0.48, 0.58,
             "And worse — on hard tasks (GPQA below MMLU): where audit is needed most, the explanation can be trusted least.",
             size=11.5, color=DEEP, line_spacing=1.12)
    gold_callout(s, 0.55, 5.32, 12.25, 0.96,
                 "What to do: a human validator checks the RESULT against an independent source (a database, document, calculation, expert), not the attached \"chain of reasoning\". Control based on the model's self-explanation is not control.",
                 size=13.5)
    speaker_notes(s, load_notes("s07"))


def build_s17(p):
    """assertion_visual (§3.3/§3.5) — FT narrowed to behavior: criteria
    "what goes where" + a hybrid is the norm. Table knowledge→RAG / behavior→PEFT /
    cheaper→distillation / deterministic→code. (Moved from s14; s14 is now
    about distillation as a separate technique.)"""
    s = blank(p)
    slide_title(s, "What here is knowledge, what is behavior, what is deterministic.", size=25)
    text_box(s, 0.55, 1.12, 12.25, 0.44,
             "The design question is not \"RAG or fine-tuning\". Split the problem along axes: knowledge → RAG, behavior → PEFT, cheaper → distillation, deterministic → code.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    ocean_box(s, 0.40, 1.74, 12.55, 3.34)
    tx, ty = 0.55, 1.86
    headers = ["If the task requires…", "→ the right tool", "…and NOT this one, because"]
    col_w = [3.95, 3.55, 4.75]
    rows = [
        ("knowledge changes / freshness, provenance needed", "RAG (or a long context for a small corpus)",
         "not fine-tuning: knowledge will go stale, retraining is expensive, risk of forgetting", False),
        ("stable behavior / tone / format / policy", "fine-tuning (PEFT)",
         "not RAG: it feeds knowledge into the context but doesn't change the model's manner", False),
        ("reduce cost / latency on a narrow task", "fine-tune a teacher + distill into a student",
         "two separate techniques in tandem; prompt/RAG don't shrink the model size", False),
        ("a deterministic, verifiable answer", "plain code, no AI",
         "neither RAG nor FT: AI adds non-determinism with no gain", True),
    ]
    hh, rh = 0.50, 0.66
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID, radius=False)
        text_box(s, cx + 0.14, ty, col_w[j] - 0.28, hh, hd,
                 size=12.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for ri, (c0, c1, c2, isgold) in enumerate(rows):
        bgrow = GOLD_TINT if isgold else (WHITE if ri % 2 == 0 else SURFACE)
        cx = tx
        for j, cc in enumerate([c0, c1, c2]):
            filled_rect(s, cx, yy, col_w[j], rh, bgrow,
                        stroke=(GOLD if isgold else SOFT_GREY),
                        stroke_pt=(1.5 if isgold else 0.75))
            col = (DEEP if isgold else MID) if j == 1 else DEEP
            text_box(s, cx + 0.14, yy + 0.04, col_w[j] - 0.28, rh - 0.08, cc,
                     size=11.5, bold=(j == 1), color=col,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 5.24, 12.25, 1.00,
                 "What to do: a hybrid is the norm WHERE the task has BOTH a knowledge problem AND a behavior problem. No behavior problem — no FT needed, even when there's RAG. Each component is added to meet its own requirement — the same ladder rule.",
                 size=13.5)
    speaker_notes(s, load_notes("s17"))


def build_s19b(p):
    """assertion_visual (§4.1) — economics: agent ×N tokens vs chat,
    multi-agent another ×15; prompt caching goes from "nice reference" → "necessity".
    Baseline — the cost of a single chat turn."""
    s = blank(p)
    slide_title(s, "An agent is not a \"pricier chat\". It's a different cost class.", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.44,
             "The comparison baseline is one chat turn (tens of cents). Each step up multiplies token spend, it doesn't just add to it.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # 3 escalating cost bars (relative to chat-turn = 1×)
    tiers = [
        ("Single LLM call (chat)", "×1", 0.11, LIGHT, "the base cost of a request"),
        ("Single agent (loop)", "≈50×", 0.36, MID, "plan→act→check→repeat: many passes per task"),
        ("Multi-agent", "another ×15", 1.0, TEAL, "on top of the agent — coordinating several agents (Anthropic, 2025)"),
    ]
    by, bx = 1.86, 0.55
    row_h = 1.14
    max_w = 8.6
    for nm, mult, frac, col, sub in tiers:
        ocean_box(s, bx, by, 12.25, row_h - 0.14)
        text_box(s, bx + 0.24, by + 0.12, 3.4, 0.40, nm, size=14.5, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, bx + 0.24, by + 0.54, 3.4, 0.40, sub, size=11,
                 italic=True, color=SLATE, line_spacing=1.05)
        filled_rect(s, bx + 3.75, by + 0.30, max_w, 0.40, SOFT_GREY,
                    radius=True, radius_adj=0.3)
        filled_rect(s, bx + 3.75, by + 0.30, max(0.5, max_w * frac), 0.40,
                    col, radius=True, radius_adj=0.3)
        chip(s, bx + 3.75 + max(0.5, max_w * frac) - 1.35, by + 0.22, 1.30,
             0.56, mult, fill=(GOLD if col is TEAL else col),
             color=(DEEP if col is TEAL else WHITE), size=15)
        by += row_h
    gold_callout(s, 0.55, 5.28, 9.35, 1.42,
                 "What to do: budget BEFORE choosing the architecture. Prompt caching (don't recompute the unchanged prefix) in an agent is no longer a \"nice reference\" but a necessity: without it the ×N loop cost becomes unaffordable.",
                 size=13.5)
    # #185: the real internet meme "Batman slap" — "an agent is just a pricier chat"
    # → a slap: it's a different cost class, not a "pricier chat".
    bmw, bmh = 1.85, 1.79
    bmx, bmy = 10.55, 5.05
    filled_rect(s, bmx - 0.05, bmy - 0.05, bmw + 0.10, bmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s19b-batman-en.png", bmx, bmy, bmw, bmh)
    footer(s, "The multipliers are order-of-magnitude (Anthropic, 2025), not an exact tariff; the absolute price depends on the model and task.")
    speaker_notes(s, load_notes("s19b"))


def build_s20(p):
    """assertion_visual (§4.1) — MCP: N×M→N+M ("USB-C"); ~11% of the catalog is actually
    runnable; the trust turn (30+ CVE/60 days, path traversal 82%). Ease of
    connection ≠ security."""
    s = blank(p)
    slide_title(s, "MCP: connecting a tool is now a matter of minutes.", size=25)
    # LEFT — N×M -> N+M
    lx, ly, lw, lh = 0.55, 1.30, 6.05, 2.30
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "cable", lx + 0.24, ly + 0.18, 0.44, "mid")
    text_runs(s, lx + 0.82, ly + 0.20, lw - 1.0, 0.42, [
        {"text": "MCP", "size": 17, "bold": True, "color": MID},
        {"text": "  — a connection standard (already covered above)", "size": 12.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 0.82, lw - 0.5, 1.30,
             "Describe a tool once as an MCP server — and every model-client sees it. Connecting became a matter of minutes. But the ease of connecting says nothing about what exactly you're connecting.",
             size=13, color=DEEP, line_spacing=1.20)
    # RIGHT — the catalog denominator
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.40,
             "Read the catalog critically", size=14, bold=True, color=DEEP)
    text_runs(s, rx + 0.24, ly + 0.62, rw - 0.48, 0.80, [
        {"text": "\"up to 90,000 servers\"  →  actually runnable ≈ ", "size": 13,
         "color": DEEP},
        {"text": "10,000", "size": 16, "bold": True, "color": MID},
        {"text": "  = ", "size": 13, "color": DEEP},
        {"text": "11%", "size": 20, "bold": True, "color": GOLD_TINT and MID},
        {"text": " of the catalog.", "size": 13, "color": DEEP},
    ], line_spacing=1.18)
    text_box(s, rx + 0.24, ly + 1.44, rw - 0.48, 0.72,
             "The rest are duplicates, broken and stubs. And even the working 10 thousand aren't checked for security.",
             size=12, color=DEEP, line_spacing=1.16)
    # BOTTOM — trust turn (security)
    ty2, th = 3.78, 1.60
    ocean_box(s, 0.55, ty2, 12.25, th, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    icon(s, "shield-alert", 0.80, ty2 + 0.22, 0.46, "teal")
    text_box(s, 1.40, ty2 + 0.22, 11.1, 0.42,
             "The trust turn: standardizing the connection ≠ security of what's connected — and it sharpens it.",
             size=15, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    facts = [
        ("30+ CVE", "in a 60-day window against MCP servers (~43% — command injection)"),
        ("82%", "path traversal among 2,614 checked implementations"),
    ]
    fx = 0.80
    for big, sub in facts:
        text_box(s, fx, ty2 + 0.78, 1.6, 0.46, big, size=22, bold=True,
                 color=TEAL)
        text_box(s, fx + 1.65, ty2 + 0.80, 4.35, 0.66, sub, size=11.5,
                 color=DEEP, line_spacing=1.12)
        fx += 6.15
    gold_callout(s, 0.55, 5.56, 12.25, 0.90,
                 "What to do: ease of connection is not an argument for connecting. Each MCP server = someone else's code in your environment + a carrier of injection into the context. You connect it to meet the task's requirement, with an assessment of the new trust boundary.",
                 size=13)
    footer(s, "Current figures on the MCP ecosystem and the CVE timeline — see sources.")
    speaker_notes(s, load_notes("s20"))

def build_s22a_multi(p):
    """assertion_visual (§4.3) — multi-agent by default is NOT an upgrade: p^n
    (95%×10≈60%); topology swarm 17.2× vs coordinator 4.4×. Anthropic:
    "works mainly because it helps spend enough tokens"."""
    s = blank(p)
    slide_title(s, "Multi-agent by default is not a win — it is a risk multiplier.", size=24)
    text_box(s, 0.55, 1.10, 12.25, 0.44,
             "In a chain of n steps, where a failure of any one ruins the result, reliabilities multiply: success ≈ pⁿ. The engineering intuition \"each step almost always works\" is mathematically false.",
             size=13, italic=True, color=MID, line_spacing=1.14)
    # LEFT — p^n visualization (bar per step count)
    lx, ly, lw, lh = 0.55, 1.72, 6.05, 3.28
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.24, ly + 0.12, lw - 0.48, 0.34,
             "Chain reliability = pⁿ", size=14, bold=True, color=DEEP)
    pns = [
        ("95% × 10 steps", 60, MID),
        ("90% × 10 steps", 35, TEAL),
    ]
    pb = ly + 0.58
    mw = lw - 2.0
    for nm, pct, col in pns:
        text_box(s, lx + 0.24, pb, lw - 0.48, 0.26, nm, size=12.5, bold=True,
                 color=DEEP)
        filled_rect(s, lx + 0.24, pb + 0.30, mw, 0.36, SOFT_GREY, radius=True,
                    radius_adj=0.3)
        filled_rect(s, lx + 0.24, pb + 0.30, mw * pct / 100.0, 0.36, col,
                    radius=True, radius_adj=0.3)
        text_box(s, lx + 0.24 + mw + 0.10, pb + 0.26, 1.3, 0.44,
                 f"~{pct}%", size=17, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        pb += 0.86
    filled_rect(s, lx + 0.24, pb + 0.04, lw - 0.48, 0.86, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.08)
    text_box(s, lx + 0.40, pb + 0.10, lw - 0.78, 0.74,
             "95% per step sounds reliable — but 0.95¹⁰ ≈ 0.60. More agents = more steps = lower overall reliability.",
             size=11.5, color=DEEP, line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
    # RIGHT — topology + Anthropic quote
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, 1.46, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.36,
             "Topology decides: coordinator > swarm", size=14, bold=True,
             color=TEAL)
    text_runs(s, rx + 0.24, ly + 0.58, rw - 0.48, 0.80, [
        {"text": "A \"swarm\" of equal agents amplifies errors ", "size": 12,
         "color": DEEP},
        {"text": "17.2×", "size": 16, "bold": True, "color": TEAL},
        {"text": ";  a single coordinator — only ", "size": 12, "color": DEEP},
        {"text": "4.4×", "size": 16, "bold": True, "color": MID},
        {"text": " (Zartis/Redis).", "size": 12, "color": DEEP},
    ], line_spacing=1.18)
    ocean_box(s, rx, ly + 1.58, rw, 1.70)
    icon(s, "message-circle", rx + 0.24, ly + 1.76, 0.42, "mid")
    text_box(s, rx + 0.24, ly + 2.24, rw - 0.48, 0.94,
             "Anthropic verbatim: \"multi-agent works mainly because it helps spend enough tokens to solve the problem\" — the win comes from token volume, not from \"coordination magic\".",
             size=12, italic=True, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 5.22, 12.25, 1.02,
                 "What to do: start with one strong agent. Multi-agent — only if the task splits into WIDELY parallel independent subtasks of high value; otherwise +15× tokens and coordination overhead will not pay off.",
                 size=13)
    speaker_notes(s, load_notes("s22a_multi"))


def build_s23b(p):
    """case_study (§4.10) — catalog of 10 CLASSES of agent failures as a table:
    class → case/date → lesson. The learning unit is the CLASS, not the case.
    Base: 188 of 344 enterprise-relevant — the agent broke it itself."""
    s = blank(p)
    slide_title(s, "Agent failures are CLASSES, not a list of oddities.", size=24)
    text_runs(s, 0.55, 1.08, 12.25, 0.44, [
        {"text": "Of ", "size": 13, "color": MID, "italic": True},
        {"text": "344", "size": 14, "bold": True, "color": MID},
        {"text": " business-significant AI incidents, in ", "size": 13,
         "color": MID, "italic": True},
        {"text": "188", "size": 14, "bold": True, "color": TEAL},
        {"text": " (≈55%) the autonomous system caused damage in production WITHOUT an attacker — the agent broke everything itself. Remember the class, not the date.",
         "size": 13, "color": MID, "italic": True},
    ], line_spacing=1.14)
    ocean_box(s, 0.40, 1.58, 12.55, 4.60)
    tx, ty = 0.52, 1.68
    headers = ["Failure class", "Case · date", "Lesson learned"]
    col_w = [3.75, 3.30, 5.25]
    rows = [
        ("Destruction + excess rights", "PocketOS · 04.2026", "System prompt ≠ control; hard boundary + least-privilege"),
        ("Injection without a click (zero-click)", "EchoLeak / M365 Copilot", "The \"lethal trifecta\" = exploitable"),
        ("Uncontrolled spend", "$48k/14h · $1.3M/30d", "No success criterion and no hard ceiling → the loop won't end"),
        ("Package hallucination", "slopsquatting · 19.7%", "An invented package name → the attacker registers it; don't trust without a registry"),
        ("Malicious MCP server", "postmark-mcp · 09.2025", "An MCP server = an unverified supply chain; version swap (rug-pull)"),
        ("Multi-agent cascade", "61% of cascades from upstream", "Where it broke ≠ where it showed; errors multiply"),
        ("Legal liability", "OLG Hamm · Air Canada", "\"It was the bot's answer\" — no defense; the output belongs to the company"),
        ("Loop without a budget", "$4,200 / 63 h", "\"Try until it works\" with no limit = literally"),
        ("Error accumulation", "reliability as pⁿ", "pⁿ, not an average; \"tune the model\" is a weak lever"),
        ("Multi-agent fragility", "Cognition, 2025", "For tasks with dependencies, multi-agent is worse than one"),
    ]
    hh = 0.40
    rh = (4.60 - 0.20 - hh) / len(rows)
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID, radius=False)
        text_box(s, cx + 0.12, ty, col_w[j] - 0.24, hh, hd,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for ri, (c0, c1, c2) in enumerate(rows):
        bgrow = WHITE if ri % 2 == 0 else SURFACE
        cx = tx
        for j, cc in enumerate([c0, c1, c2]):
            filled_rect(s, cx, yy, col_w[j], rh, bgrow, stroke=SOFT_GREY,
                        stroke_pt=0.5)
            text_box(s, cx + 0.12, yy, col_w[j] - 0.24, rh, cc,
                     size=10.5, bold=(j == 0), color=(DEEP if j == 0 else DEEP),
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 6.30, 12.25, 0.78,
                 "What to do: faced with a new incident — name ITS CLASS. The class sets the countermeasure (limit / least-privilege / check between steps), and the specific dates and amounts are not needed.",
                 size=13)
    speaker_notes(s, load_notes("s23b"))


def build_s23c(p):
    """case_study (§4.10) — deep-dive of classes with a base: PocketOS 9 sec,
    runaway $48k/$1.3M, slopsquatting 19.7% (1 in 5), cascade 61% from upstream.
    Each — a lesson + the correct alternative."""
    s = blank(p)
    slide_title(s, "Four classes up close — with a comparison base.", size=25)
    cards = [
        ("bomb", "PocketOS — a production DB in 9 seconds",
         "The agent found a token with unlimited rights in someone else's file and deleted the volume + all backups. The nearest recoverable one was 3 months old.",
         "Lesson: the system prompt is not a security control; you need a hard boundary."),
        ("flame", "Uncontrolled spend: $48k / 14h · $1.3M / 30d",
         "A request with no success criterion — the planner expanded and never finished. Base: a normal session is cents-to-dollars; this is 3–5 orders of magnitude higher.",
         "Lesson: a completion criterion + a hard ceiling + an emergency stop are mandatory."),
        ("package", "Slopsquatting — 1 in 5 imports",
         "Over 576,000 code samples, 19.7% of package references are hallucinations. The attacker registers the invented name → the next agent installs someone else's code.",
         "Lesson: pin versions/hashes, lock files, review of new dependencies."),
        ("git-merge", "Cascade — 61% of errors from upstream",
         "Across 73 incidents: in 61% the root was in an upstream layer (retrieval/plan), not where the failure became visible. The next agent takes the error as fact.",
         "Lesson: checks between steps, fewer hops, traceability."),
    ]
    gx, gy = 0.55, 1.28
    cw, chh = 6.05, 2.42
    gap = 0.15
    for i, (ic, title, body, lesson) in enumerate(cards):
        x = gx + (i % 2) * (cw + gap)
        y = gy + (i // 2) * (chh + gap)
        ocean_box(s, x, y, cw, chh)
        filled_rect(s, x + 0.22, y + 0.22, 0.52, 0.52, TEAL, radius=True,
                    radius_adj=0.18)
        icon(s, ic, x + 0.28, y + 0.28, 0.40, "white")
        text_box(s, x + 0.86, y + 0.22, cw - 1.06, 0.52, title, size=14,
                 bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        text_box(s, x + 0.24, y + 0.86, cw - 0.48, 1.02, body, size=11.5,
                 color=DEEP, line_spacing=1.14)
        filled_rect(s, x + 0.24, y + chh - 0.50, cw - 0.48, 0.40, GOLD_TINT,
                    stroke=GOLD, stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(s, x + 0.36, y + chh - 0.50, cw - 0.72, 0.40, lesson,
                 size=10.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    gold_callout(s, 0.55, 6.42, 12.25, 0.72,
                 "What to do: in all four the root is the same — the agent was applied without an external boundary (of rights, budget, validation). The boundary is placed OUTSIDE the agent; a prompt cannot replace it.",
                 size=13)
    speaker_notes(s, load_notes("s23c"))


def build_s24(p):
    """assertion_visual (§4.8) — a per-feature data map: ZDR does NOT cover
    third-party/MCP/Files/batch; NYT v OpenAI litigation hold outlived "30 days".
    The more agentic — the more data outside ZDR."""
    s = blank(p)
    slide_title(s, "\"We have ZDR\" ≠ \"data is protected across the whole chain\".", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.44,
             "ZDR (zero data retention) covers the core model calls — but not the whole architecture. An agent = a model + tools, and tools are often outside ZDR.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # LEFT — what ZDR does NOT cover
    lx, ly, lw, lh = 0.55, 1.82, 6.05, 3.10
    ocean_box(s, lx, ly, lw, lh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "circle-slash", lx + 0.24, ly + 0.18, 0.44, "teal")
    text_box(s, lx + 0.82, ly + 0.20, lw - 1.0, 0.42,
             "ZDR does NOT cover", size=15, bold=True, color=TEAL,
             anchor=MSO_ANCHOR.MIDDLE)
    outs = ["Files API", "batch processing", "code execution in containers",
            "MCP connector", "third-party integrations",
            "consumer plans"]
    oy = ly + 0.76
    for i, o in enumerate(outs):
        col_i = i % 2
        ox = lx + 0.28 + col_i * 2.90
        oyy = oy + (i // 2) * 0.56
        circle(s, ox, oyy + 0.06, 0.11, TEAL)
        text_box(s, ox + 0.22, oyy, 2.60, 0.50, o, size=12, color=DEEP,
                 line_spacing=1.02, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.28, ly + 2.54, lw - 0.56, 0.46,
             "The more agentic the architecture — the more of its data passes through these links.",
             size=12, bold=True, color=TEAL, line_spacing=1.10)
    # RIGHT — litigation hold overrides policy
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "gavel", rx + 0.24, ly + 0.18, 0.44, "gold")
    text_box(s, rx + 0.82, ly + 0.20, rw - 1.0, 0.42,
             "A court order outlives your policy", size=14, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, rx + 0.28, ly + 0.80, rw - 0.56, 1.30,
             "NYT v. OpenAI (2025): the court ordered ALL ChatGPT logs preserved as evidence. The contractual \"30 days\" policy was powerless against a litigation hold from someone else's dispute.",
             size=13, color=DEEP, line_spacing=1.20)
    text_box(s, rx + 0.28, ly + 2.14, rw - 0.56, 0.86,
             "Data that has left your perimeter lives by rules you do not control.",
             size=13, bold=True, color=MID, line_spacing=1.16)
    gold_callout(s, 0.55, 5.06, 9.05, 1.62,
                 "What to do: map the data per feature before production — what data, through which link, with what retention policy, which links are third-party. Regulated/sensitive data — only with ZDR/BAA or locally (on-prem).",
                 size=13)
    # #185: real internet meme "Always has been" — "data outside ZDR?" / "always
    # has been": what leaves the perimeter is outside your retention policy.
    amw, amh = 3.00, 1.69
    amx, amy = 9.80, 5.06
    filled_rect(s, amx - 0.05, amy - 0.05, amw + 0.10, amh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.04)
    add_image(s, WEB / "s24-alwayshasbeen-en.png", amx, amy, amw, amh)
    speaker_notes(s, load_notes("s24"))


def build_s28(p):
    """summary (§5.3) — summary table "mechanism → boundary → what to do" across
    all architectures (L2 s38 pattern). ≤2 lines/cell, NO callback to s01.
    The bottom row — "not AI at all". Gold takeaway."""
    s = blank(p)
    slide_title(s, "Takeaway: mechanism → its boundary → what to do.", size=25)
    ocean_box(s, 0.40, 1.20, 12.55, 5.06)
    tx, ty = 0.52, 1.30
    headers = ["Mechanism", "Where it breaks (boundary)", "What to do"]
    col_w = [3.05, 4.55, 4.70]
    rows = [
        ("Prompt / role", "a role-persona changes tone, not accuracy", "accuracy — via context and RAG, not via \"you are an expert\""),
        ("Chain-of-thought", "faithfulness is low (~25–39%)", "check the result, not the self-explanation"),
        ("RAG", "\"retrieved\" ≠ \"retrieved correctly\"", "grounding in a source + retrieval metrics + \"I don't know\" as the norm"),
        ("Fine-tuning / PEFT", "changes behavior, not knowledge; forgetting", "PEFT + a check loop + versioning; knowledge → RAG"),
        ("Agent loop", "plan→act→check→repeat — 4 points of failure", "check against an external criterion; a hard ceiling on the loop"),
        ("Equipment / memory", "each slot is a trade-off, not a win", "add a slot on demand, with a check"),
        ("Multi-agent", "pⁿ: 95%×10 ≈ 60%; +15× tokens", "one strong agent by default"),
        ("Security", "injection × broad rights; ZDR isn't everything", "least-privilege + a per-feature data map"),
        ("\"Not AI at all\"", "deterministic + verifiable", "plain code — cheaper, more predictable, auditable", True),
    ]
    hh = 0.46
    rh = (5.06 - 0.20 - hh) / len(rows)
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID, radius=False)
        text_box(s, cx + 0.12, ty, col_w[j] - 0.24, hh, hd,
                 size=12.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for ri, row in enumerate(rows):
        isgold = len(row) == 4 and row[3]
        c0, c1, c2 = row[0], row[1], row[2]
        bgrow = GOLD_TINT if isgold else (WHITE if ri % 2 == 0 else SURFACE)
        cx = tx
        for j, cc in enumerate([c0, c1, c2]):
            filled_rect(s, cx, yy, col_w[j], rh, bgrow,
                        stroke=(GOLD if isgold else SOFT_GREY),
                        stroke_pt=(1.5 if isgold else 0.5))
            text_box(s, cx + 0.12, yy, col_w[j] - 0.24, rh, cc,
                     size=11, bold=(j == 0 or isgold), color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 6.36, 12.25, 0.72,
                 "To know a tool is to know its boundaries. Choosing an architecture = find the lowest rung that meets the task's requirement.",
                 size=13.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s28"))


def build_s31(p):
    """qa_minimal — dedicated final Q&A slide (#239, Lecture 1 s31 style).

    A large "Q&A" 120pt centered in DEEP; "Thank you" 36pt below; the lecturer's
    contact details small in the bottom-right corner (filled in before the
    lecture). White background, no footer and no roadmap-bar."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    # #185: real internet meme "Waiting Skeleton" — "waiting for your questions".
    # Left column; the Q&A text is shifted right so the meme does not overlap.
    kmx, kmy, kmw, kmh = 0.85, 1.55, 3.25, 4.40
    ocean_box(s, kmx, kmy, kmw, kmh)
    add_image(s, WEB / "s31-skeleton-en.png", kmx + 0.14, kmy + 0.14,
              kmw - 0.28, kmh - 0.28)
    text_box(s, x=4.35, y=2.05, w=8.45, h=2.30, text="Q&A",
             size=120, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    text_box(s, x=4.35, y=4.55, w=8.45, h=0.78,
             text="Thank you", size=36, bold=False, color=MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.2)
    # The instructor's contacts are NOT on the visible layer (scaffold) — a VFY
    # plate is in the speaker notes (filled in by the owner before the lecture at GATE B).
    speaker_notes(s, load_notes("s31"))


# ============================================================
# Main
# ============================================================
def main():
    # U-9: 2-part deck spec — loader reads deck.yaml + deck-part2.yaml,
    # validates 36-slide order + cascade lock (s01–s30 not renumbered).
    spec = load_deck()
    if spec is not None:
        n = len(spec["slides"])
        print(f"deck spec OK — {n} slides (deck.yaml + deck-part2.yaml), "
              f"version {spec['deck'].get('version')}, "
              f"totals {spec['totals'].get('slides')}")

    p = setup_pres()
    # v5 (issue #185) presentation order — 51 slides (authoritative ordered
    # list per deck-v5-inventory-draft §3; §9's "54" double-counted dividers).
    #   R0: s01(meme) s02 s02a s03 s04   (s01b removed — #185/#312)
    #   R1: s04a(div) s05 s05a s05c s05b s06 s07 s08 s08a
    #   R2: s09(div) s10 s11 s12 s13
    #   R3: s13a(div) s13b s15 s17 s14 s16
    #   R4: s18(div) s19 s19b s20 s21 s22 s22a_multi s22b s22c s22d
    #        s22e s25 s24 s25b s23 s23b s23c
    #   R5: s25a(div) s26 s27 s27b s28 s29 s30 s31
    # v5b (issue #185 WP8): +5 "classical baseline" slides — one per
    # section §1–§5, right AFTER the section divider, BEFORE the AI part. 51→56.
    builders = [
        # R0 — Opening (5) — s01b removed (#185/#312): the hook is already on s01,
        # Air Canada remains a §2 case (s13).
        build_s01, build_s02, build_s02a, build_s03, build_s04,
        # R1 — Prompt (div + classic-base + 8)
        build_s04a, build_s_classic_prompt, build_s05, build_s05a, build_s05c,
        build_s05b, build_s06, build_s07, build_s08, build_s08a,
        # R2 — RAG (div + classic-base + 4)
        build_s09, build_s_classic_rag, build_s10, build_s11, build_s12, build_s13,
        # R3 — Fine-tune (div + classic-base + 5)
        build_s13a, build_s_classic_ft, build_s13b, build_s15, build_s17,
        build_s14, build_s16,
        # R4 — Agents (div + classic-base + 16)
        build_s18, build_s_classic_agents, build_s19, build_s19b, build_s20,
        build_s21, build_s22, build_s22a_multi, build_s22b, build_s22c,
        build_s22d, build_s22e, build_s25, build_s24, build_s25b, build_s23,
        build_s23b, build_s23c,
        # R5 — Framework (div + classic-base + 7)
        build_s25a, build_s_classic_framework, build_s26, build_s27, build_s27b,
        build_s28, build_s29, build_s30, build_s31,
    ]
    # sid list — MUST match `builders` order 1:1 (display order, 55 slides).
    sids = [
        "s01", "s02", "s02a", "s03", "s04",
        "s04a", "s-classic-prompt", "s05", "s05a", "s05c", "s05b", "s06",
        "s07", "s08", "s08a",
        "s09", "s-classic-rag", "s10", "s11", "s12", "s13",
        "s13a", "s-classic-ft", "s13b", "s15", "s17", "s14", "s16",
        "s18", "s-classic-agents", "s19", "s19b", "s20", "s21", "s22",
        "s22a_multi", "s22b", "s22c", "s22d", "s22e", "s25", "s24", "s25b",
        "s23", "s23b", "s23c",
        "s25a", "s-classic-framework", "s26", "s27", "s27b", "s28", "s29",
        "s30", "s31",
    ]
    assert len(builders) == 55, f"expected 55 builders, got {len(builders)}"
    assert len(sids) == 55, f"expected 55 sids, got {len(sids)}"

    total = len(builders)
    inject_report = {}
    for idx, (b, sid) in enumerate(zip(builders, sids)):
        b(p)
        slide = p.slides[idx]
        # (1) inject small superscript [N] markers at claim anchors
        inject_report[sid] = R.inject_ref_markers(slide, sid)
        # (2) bottom clickable numbered source list (fold any footer caveat in)
        if sid in R.SLIDE_REFS:
            ftext = _FOOTER_TEXT.get(id(slide))
            fshape = _FOOTER_TEXT.get("_shapes", {}).get(id(slide))
            if fshape is not None:
                # remove the standalone footer; its words survive as ref tail
                fshape._element.getparent().remove(fshape._element)
            R.refs_of_slide(slide, sid, y=7.02, tail=ftext)
        # (3) speaker notes already carry the "Sources:" block + [N] markers
        # (baked into slides/*.md by patch_notes.py — single source of truth,
        # so slide-[N] and notes-[N] do not diverge; builder's speaker_notes(
        # load_notes(sid)) picks it up). Nothing to do here.
        # (4) muted page number "N / 40" bottom-right on every slide
        R.page_number(slide, idx + 1, total)

    # verification print — any anchor that failed to match
    missed = [(sid, a) for sid, rep in inject_report.items()
              for (a, ok) in rep if not ok]
    if missed:
        print("!! UNMATCHED ANCHORS:")
        for sid, a in missed:
            print(f"   {sid}: {a[:70]}")
    else:
        print("all ref anchors matched OK")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"saved {OUT} — {len(p.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
