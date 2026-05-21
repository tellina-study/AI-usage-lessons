"""
Full 39-slide build of Лекции 11 «AI в дискретном и процессном производстве».

Source-of-truth: deck.yaml v1 + chapter v5 multi-part (~28k слов) + slides/*.md.

Issue #127 · downstream от chapter (book-first).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml ≈ 75 мин.

Lec-N-1 pattern compliance: match lec-09 (cover + lecture-map + 5 section dividers +
dedicated Q&A; top progress bar только на dividers + cover).

Build via: python3 build_lec11.py — generates lec-11.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
RED_WARN = RGBColor(0xC0, 0x39, 0x2B)
ROADMAP = RGBColor(0xD9, 0xE2, 0xEC)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "lec-11.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *, size=16, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except: pass
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline_box(slide, x, y, w, h, lines, *, size=14, bold=False, color=DEEP,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
                  line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except: pass
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w, h):
    if not path or not Path(path).exists():
        return None
    try:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                         width=Inches(w), height=Inches(h))
    except Exception as e:
        print(f"Image fail {path}: {e}")
        return None


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def roadmap_bar(slide, current_section):
    """5-section roadmap bar at top of section dividers + cover."""
    sections = ["Общее", "Дискретное", "Процессное", "Рамка", "Замыкание"]
    bar_y = 0.4
    bar_h = 0.32
    total_w = 12.33
    seg_w = total_w / 5
    for i, name in enumerate(sections):
        x = 0.5 + i * seg_w
        is_active = (i == current_section)
        rectangle(slide, x, bar_y, seg_w - 0.05, bar_h,
                  fill=GOLD if is_active else ROADMAP)
        text_box(slide, x, bar_y + 0.04, seg_w - 0.05, bar_h - 0.04,
                 f"{i+1}. {name}", size=10, bold=is_active,
                 color=DEEP if is_active else SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, text, *, y=7.05):
    text_box(slide, 0.5, y, 12.33, 0.35, text,
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)


def attribution(slide, text, *, x=0.5, y=6.95, w=12.33):
    text_box(slide, x, y, w, 0.3, text,
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)


# ========== SECTION 0 ==========

def s01_hero_tesla(p):
    """s01 — hero hook: Tesla Giga Press BEFORE/AFTER."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    # Hero image left (50% area)
    img_path = ASSETS / "screenshots" / "s01-tesla-giga-press.png"
    if img_path.exists():
        add_image(slide, img_path, 0.5, 0.5, 6.5, 4.8)
    else:
        rectangle(slide, 0.5, 0.5, 6.5, 4.8, fill=SOFT_GREY)
        text_box(slide, 0.5, 2.5, 6.5, 0.6, "[Tesla Giga Press hero]",
                 size=16, color=SLATE, align=PP_ALIGN.CENTER)
    # Caption under hero
    attribution(slide, "Tesla Giga Press · Idra OL 6100 CS · Fremont, 2020 · Wikimedia CC-BY-SA",
                x=0.5, y=5.4, w=6.5)
    # Title block right
    multiline_box(slide, 7.3, 0.7, 5.7, 5.0, [
        ("Tesla отступила дважды.", {"size": 30, "bold": True, "color": DEEP}),
        ("Компании не учатся один раз.", {"size": 24, "bold": True, "color": MID}),
        ("", {"size": 10}),
        ("Май 2024:", {"size": 14, "bold": True, "color": GOLD}),
        ("Tesla отказалась от next-generation gigacasting для Model 2.", {"size": 14, "color": DEEP}),
        ("", {"size": 10}),
        ("Апрель 2018:", {"size": 14, "bold": True, "color": MID}),
        ("«Excessive automation at Tesla was a mistake. To be precise, my mistake. Humans are underrated.»", {"size": 13, "italic": True, "color": DEEP}),
        ("— Илон Маск, X, 13.04.2018", {"size": 11, "italic": True, "color": SLATE}),
    ], line_spacing=1.2)
    # Central question
    rounded_box(slide, 0.5, 5.9, 12.33, 1.0, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 0.8, 6.05, 11.8, 0.8,
             "MIT Sloan 2025: 95% AI-пилотов не доходят до production. Где AI работает, где нет — и как решать инженеру?",
             size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "В мае 2024 Tesla отступила от next-generation gigacasting. Это вторая отмена — первая была в апреле 2018, в кризис Model 3, когда Маск признал чрезмерную автоматизацию своей ошибкой. Компании не учатся один раз. Параллельно GE сожгла 4 миллиарда на Predix, IBM Watson продан за parts, Foxconn Wisconsin превратился в дата-центр Microsoft. И поверх — MIT Sloan: 95% пилотов не доходят до production. Это лекция про то, где AI работает, где не работает, и как решать инженеру, в какой колонне он стоит — дискретное или процессное.")


def s02_cover(p):
    """s02 — cover."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=0)
    # Large decorative «11»
    text_box(slide, 0.5, 1.2, 4.0, 4.5, "11",
             size=240, bold=True, color=ROADMAP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    # Title
    multiline_box(slide, 4.5, 1.5, 8.3, 2.0, [
        ("Лекция 11", {"size": 22, "bold": True, "color": LIGHT}),
        ("AI в дискретном и процессном", {"size": 38, "bold": True, "color": DEEP}),
        ("производстве", {"size": 38, "bold": True, "color": DEEP}),
    ], line_spacing=1.05)
    # Meta
    multiline_box(slide, 4.5, 4.3, 8.3, 1.0, [
        ("Модуль 2 · 75 минут + Q&A", {"size": 16, "color": LIGHT}),
        ("Студенты-инженеры 3 курса", {"size": 16, "color": LIGHT}),
    ])
    # LO summary in rounded box
    rounded_box(slide, 4.5, 5.4, 8.3, 1.4)
    multiline_box(slide, 4.7, 5.55, 8.0, 1.2, [
        ("Цели лекции:", {"size": 14, "bold": True, "color": MID}),
        ("LO1 — назвать инструменты дискретного и процессного AI · LO2 — критически оценить vendor-claim ·", {"size": 12, "color": DEEP}),
        ("LO7 — различить chat-помощник vs autonomous controller · LO8 — сформулировать «когда AI не нужен»", {"size": 12, "color": DEEP}),
    ], line_spacing=1.2)
    footer(slide, "Курс «Применение AI в инженерии» · 2026")
    add_notes(slide, "Лекция 11, первая из двух про реальное производство. Сегодня дискретное и процессное. Лекция 12 продолжит с фокусом на цифровые двойники. Пять целей по таксономии Блума: назвать инструменты, оценить vendor-claim, различить архитектурные классы AI, сформулировать когда AI не подходит.")


def s03_lecture_map(p):
    """s03 — 5 horizontal section cards."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Маршрут лекции — пять разделов",
             size=28, bold=True, color=DEEP)
    sections = [
        ("1", "Общее", "12 мин", "Adoption · OT/IT · foundation models · трио hype-collapse", LIGHT),
        ("2", "Дискретное", "17 мин", "CV-инспекция · разметка · PdM · коботы · Tesla 2018", MID),
        ("3", "Процессное", "17 мин", "Мягкие сенсоры · MPC/RL · регуляторика · РФ", TEAL),
        ("4", "Рамка решения", "12 мин", "4 категории критериев · альтернативы · 5-step · Pfizer", GOLD),
        ("5", "Замыкание", "6 мин", "Recap · 5 вендор-вопросов · bridge к Лекции 12", DEEP),
    ]
    card_w = 2.42
    gap = 0.05
    x0 = 0.5
    y = 1.4
    for i, (num, title, dur, desc, accent) in enumerate(sections):
        x = x0 + i * (card_w + gap)
        rounded_box(slide, x, y, card_w, 5.3)
        # Number circle
        circle(slide, x + card_w/2 - 0.5, y + 0.3, 1.0, 1.0, fill=accent)
        text_box(slide, x, y + 0.3, card_w, 1.0, num,
                 size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Title
        text_box(slide, x + 0.1, y + 1.6, card_w - 0.2, 0.6, title,
                 size=18, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        # Duration
        text_box(slide, x + 0.1, y + 2.3, card_w - 0.2, 0.4, dur,
                 size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
        # Description
        text_box(slide, x + 0.2, y + 2.85, card_w - 0.4, 2.3, desc,
                 size=11, color=DARK_GREY, align=PP_ALIGN.CENTER, line_spacing=1.3)
    footer(slide, "Раздел 4 — PAYOFF лекции · применимый инструмент для кармана")
    add_notes(slide, "Маршрут пять разделов. Сначала общее — adoption, OT/IT раскол, foundation models, трио хайп-провалов. Потом 17 минут дискретное — CV, разметка, PdM, коботы, Tesla 2018. Потом 17 минут процессное — soft sensors, MPC/RL, регуляторика, российский контекст. Потом 12 минут рамка решения — 4 категории критериев, альтернативы, worked example Pfizer Vox, 5-step framework. Раздел 4 — payoff лекции. И 6 минут на замыкание. Bridge к Лекции 12 — цифровые двойники.")


def s04_glossary(p):
    """s04 — 6 must-know terms × 2 columns."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Шесть терминов из мира промышленных систем",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Must-know для лекции · остальные acronyms — inline gloss при первом упоминании",
             size=14, italic=True, color=LIGHT)
    terms = [
        ("ISA-95", "5-уровневая иерархия систем управления (L0 процесс ↔ L4 ERP). Опорная сетка."),
        ("MES", "Manufacturing Execution System. Уровень L3 — производственные заказы, OEE-отчёты, traceability."),
        ("SCADA", "Supervisory Control And Data Acquisition. Уровень L2 — мнемосхема, диспетчерский контроль."),
        ("PLC", "Programmable Logic Controller. Уровень L1. Цикл 1–10 мс, детерминированный. Сюда AI не идёт."),
        ("OEE", "Overall Equipment Effectiveness = доступность × производительность × качество. Central метрика."),
        ("Мягкий сенсор", "Программная модель: оценивает трудно-измеряемые параметры по легко-измеряемым. Cornerstone процессного."),
    ]
    col_w = 5.9
    item_h = 1.5
    for i, (term, defn) in enumerate(terms):
        col = i % 2
        row = i // 2
        x = 0.5 + col * (col_w + 0.4)
        y = 1.85 + row * (item_h + 0.15)
        rounded_box(slide, x, y, col_w, item_h)
        # Term left strip in MID color
        rectangle(slide, x + 0.15, y + 0.15, 1.6, item_h - 0.3, fill=MID)
        text_box(slide, x + 0.15, y + 0.15, 1.6, item_h - 0.3, term,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Definition right
        text_box(slide, x + 1.95, y + 0.2, col_w - 2.1, item_h - 0.4, defn,
                 size=11, color=DEEP, line_spacing=1.3,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, "OEE — самая важная метрика лекции · возвращаемся в разделах 1, 2, 4")
    add_notes(slide, "Шесть терминов промышленных систем. ISA-95 — пятиуровневая модель Purdue. MES, SCADA, PLC — уровни L3, L2, L1 этой иерархии. OEE — общая эффективность оборудования, central метрика. Мягкий сенсор — фундаментальная концепция процессного производства, программная замена лабораторной пробы. AI заходит на L3 и L4 спокойно, на L2 со скрипом, на L1 почти никогда, на L0 никогда.")


def s05_keystone(p):
    """s05 — KEYSTONE: 2 columns + universal belt."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.8,
             "Две модели производства. AI входит в обе — но по-разному.",
             size=26, bold=True, color=DEEP)
    # Two columns
    col_w = 6.0
    col_h = 4.4
    col_y = 1.4
    # LEFT: Discrete
    rounded_box(slide, 0.5, col_y, col_w, col_h, fill=SURFACE, stroke=MID, stroke_w=2.0)
    text_box(slide, 0.65, col_y + 0.15, col_w - 0.3, 0.5, "ДИСКРЕТНОЕ",
             size=22, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(slide, 0.65, col_y + 0.75, col_w - 0.3, 0.4,
             "Единицы продукции считаются штучно",
             size=13, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    multiline_box(slide, 0.85, col_y + 1.3, col_w - 0.4, 2.0, [
        ("Авто · электроника · авионика · бытовая техника", {"size": 12, "italic": True, "color": LIGHT}),
        ("", {"size": 6}),
        ("Инструменты AI:", {"size": 13, "bold": True, "color": DEEP}),
        ("· Компьютерное зрение для контроля качества", {"size": 13, "color": DEEP}),
        ("· Коботы и worker-augmentation", {"size": 13, "color": DEEP}),
        ("· Generative process planning", {"size": 13, "color": DEEP}),
    ], line_spacing=1.3)
    # Failure mark
    rounded_box(slide, 0.65, col_y + col_h - 0.85, col_w - 0.3, 0.7,
                fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 0.85, col_y + col_h - 0.78, col_w - 0.4, 0.55,
             "Канонический провал: Tesla 2018 — «excessive automation a mistake»",
             size=11, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    # RIGHT: Process
    rounded_box(slide, 6.83, col_y, col_w, col_h, fill=SURFACE, stroke=TEAL, stroke_w=2.0)
    text_box(slide, 6.98, col_y + 0.15, col_w - 0.3, 0.5, "ПРОЦЕССНОЕ",
             size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(slide, 6.98, col_y + 0.75, col_w - 0.3, 0.4,
             "Материал течёт непрерывным потоком",
             size=13, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    multiline_box(slide, 7.18, col_y + 1.3, col_w - 0.4, 2.0, [
        ("Химия · фарма · металлургия · цемент · пивоварение", {"size": 12, "italic": True, "color": LIGHT}),
        ("", {"size": 6}),
        ("Инструменты AI:", {"size": 13, "bold": True, "color": DEEP}),
        ("· Мягкие сенсоры (real-time оценка качества)", {"size": 13, "color": DEEP}),
        ("· MPC / RL гибрид + CIRL", {"size": 13, "color": DEEP}),
        ("· Прогностическое обслуживание + edge AI", {"size": 13, "color": DEEP}),
    ], line_spacing=1.3)
    # Failure mark
    rounded_box(slide, 6.98, col_y + col_h - 0.85, col_w - 0.3, 0.7,
                fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 7.18, col_y + col_h - 0.78, col_w - 0.4, 0.55,
             "Канонический провал: F-35 ALIS — $44K/час, заменён ODIN",
             size=11, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    # Universal belt
    rounded_box(slide, 0.5, 6.05, 12.33, 0.9, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    text_box(slide, 0.7, 6.15, 12.0, 0.4,
             "ОБЩЕЕ ДЛЯ ОБЕИХ КОЛОНН",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(slide, 0.7, 6.5, 12.0, 0.4,
             "78% используют AI · 5,5% high performers · 95% пилотов не доходят до production",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    add_notes(slide, "Keystone слайд. Производство — две принципиально разные ветви. Дискретное — штучные продукты, инструменты CV / коботы / process planning, провал — Tesla 2018. Процессное — непрерывный поток, инструменты soft sensors / MPC-RL / PdM, провал — F-35 ALIS. Это разные физики, разные регуляторные карты, разные культуры. И поверх обеих — universal: 78 процентов используют AI, только 5,5 процентов high performers; 95 процентов пилотов не доходят до production. Это структурная картина 2025-2026, не аномалия.")


def s06_section1_divider(p):
    section_divider(p, 1, "Что общее для обеих моделей",
                    "Adoption · OT/IT раскол · foundation models · трио hype-collapse",
                    "Раздел 1 · 12 мин · 6 слайдов")


def section_divider(p, num, title, subtitle, meta):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=num - 1)
    # Large background number
    text_box(slide, 0.5, 1.0, 6.0, 5.5, str(num),
             size=440, bold=True, color=ROADMAP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title right
    multiline_box(slide, 5.5, 2.5, 7.5, 3.5, [
        ("РАЗДЕЛ", {"size": 18, "bold": True, "color": LIGHT}),
        (title, {"size": 38, "bold": True, "color": DEEP}),
        ("", {"size": 12}),
        (subtitle, {"size": 16, "color": MID}),
        ("", {"size": 8}),
        (meta, {"size": 13, "italic": True, "color": GOLD}),
    ], line_spacing=1.15)
    add_notes(slide, f"Раздел {num}. {title}. {subtitle}. {meta}.")


def s07_adoption(p):
    """s07 — adoption landscape with donut + bar."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "78% используют. 5,5% извлекают ценность.",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "McKinsey State of AI 2025 — adoption-value gap в ~14×",
             size=14, italic=True, color=LIGHT)
    # Hero number left
    rounded_box(slide, 0.5, 1.8, 4.0, 3.2, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    text_box(slide, 0.5, 1.95, 4.0, 1.0, "5,5%",
             size=100, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 3.3, 3.6, 1.6, [
        ("high performers", {"size": 16, "bold": True, "color": DEEP}),
        ("EBIT impact >5% от AI", {"size": 13, "italic": True, "color": DARK_GREY}),
        ("", {"size": 6}),
        ("McKinsey State of AI 2025", {"size": 11, "italic": True, "color": SLATE}),
    ], align=PP_ALIGN.CENTER, line_spacing=1.3)
    # Bar chart right
    chart_path = ASSETS / "charts" / "s07-pilot-failure.png"
    if chart_path.exists():
        rounded_box(slide, 4.7, 1.8, 8.1, 3.2)
        add_image(slide, chart_path, 4.9, 1.95, 7.7, 2.9)
    # Other sources
    multiline_box(slide, 0.5, 5.15, 12.33, 1.5, [
        ("Параллельные источники:", {"size": 14, "bold": True, "color": MID}),
        ("· MIT Sloan 2025: 95% GenAI-пилотов не доходят до production. Среднее время pilot→shutdown — 14 мес.", {"size": 13, "color": DEEP}),
        ("· RAND 2025: 80,3% AI-проектов без business value. $547B из $684B AI-инвестиций — без эффекта.", {"size": 13, "color": DEEP}),
        ("· Deloitte 2025: 42% компаний прекратили ≥1 AI-инициативу. Sunk cost per abandoned — $7,2M.", {"size": 13, "color": DEEP}),
    ], line_spacing=1.35)
    attribution(slide, "Источники: McKinsey 2025, MIT Sloan 2025, RAND 2025, Deloitte 2025 · [VFY-day-of]")
    add_notes(slide, "Запомните две цифры. 78 процентов используют AI, только 5,5 процентов — high performers. Разрыв 14 раз. Не верите McKinsey — посмотрите MIT Sloan: 95 процентов пилотов не доходят до production. RAND: 80 процентов проектов без business value, 547 миллиардов из 684 без эффекта. Deloitte: 42 процента компаний прекратили хотя бы одну AI-инициативу. Запомните паттерн: застревание на пилотной стадии — норма, а не аномалия. И OEE callback — vendor обещает –25 процентов downtime, спросите про OEE breakdown.")


def s08_market_estimates(p):
    """s08 — three analysts disagree by 5×."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Рыночные оценки 2025 расходятся в 5×",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Три аналитика, один сегмент «AI in manufacturing», три радикально разные цифры",
             size=14, italic=True, color=LIGHT)
    # Chart left
    chart_path = ASSETS / "charts" / "s08-market-divergence.png"
    if chart_path.exists():
        rounded_box(slide, 0.5, 1.8, 7.5, 4.0)
        add_image(slide, chart_path, 0.7, 1.95, 7.1, 3.7)
    # Pedagogical right
    rounded_box(slide, 8.2, 1.8, 4.6, 4.0, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    multiline_box(slide, 8.4, 2.0, 4.2, 3.8, [
        ("Расхождение 4,5×", {"size": 18, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("", {"size": 8}),
        ("Это не ошибка одного аналитика — это особенность жанра.", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Определения сегмента, methodology, scope радикально различаются между источниками.", {"size": 12, "color": DEEP}),
        ("", {"size": 8}),
        ("УРОК ДЛЯ ИНЖЕНЕРА:", {"size": 12, "bold": True, "color": MID}),
        ("читайте methodology, не верьте одной цифре.", {"size": 12, "color": DEEP}),
    ], line_spacing=1.35)
    # Gartner footer note
    multiline_box(slide, 0.5, 6.0, 12.33, 0.7, [
        ("Gartner оценивает worldwide GenAI spending в $644B на 2025 — но это все индустрии вместе, не только производство.", {"size": 12, "italic": True, "color": DARK_GREY}),
    ])
    attribution(slide, "Markets and Markets · Precedence · Fortune Business Insights · Gartner 2025 · [VFY-day-of]")
    add_notes(slide, "Первая красная тряпка. Три респектабельных аналитика дают радикально разные оценки одного сегмента в одном году. Markets and Markets — 34 миллиарда. Precedence — 8,57. Fortune — 7,6. Разница в 4,5 раза. Кто прав — никто. Они правы каждый внутри своей methodology. Markets and Markets включает сервисы, оборудование, ПО, консалтинг. Fortune — только программные платформы. Когда вендор приходит с цифрой — спрашивайте methodology. Gartner оценивает worldwide GenAI в 644 миллиарда — но это все индустрии.")


def s09_ot_it_split(p):
    """s09 — OT vs IT split-comparison."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "OT/IT раскол — фундаментальный structural divide",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "AI приходит из IT в OT и упирается в этот раскол",
             size=14, italic=True, color=LIGHT)
    # Two columns
    col_y = 1.8
    col_h = 4.5
    col_w = 6.0
    # OT left
    rounded_box(slide, 0.5, col_y, col_w, col_h, fill=SURFACE, stroke=MID, stroke_w=2.0)
    rectangle(slide, 0.5, col_y, col_w, 0.6, fill=MID)
    text_box(slide, 0.5, col_y + 0.05, col_w, 0.5, "OT · Operational Technology",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.75, col_y + 0.8, col_w - 0.5, 3.5, [
        ("Где живёт:", {"size": 11, "bold": True, "color": LIGHT}),
        ("PLC · SCADA · DCS · конвейерные приводы · реакторы", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Цикл управления:", {"size": 11, "bold": True, "color": LIGHT}),
        ("1–10 мс детерминированный", {"size": 13, "bold": True, "color": GOLD}),
        ("", {"size": 6}),
        ("Консистентность: strong", {"size": 12, "color": DEEP}),
        ("Audit trail: обязательный, регуляторный", {"size": 12, "color": DEEP}),
        ("Sertification: SIL 2/3 для safety-critical", {"size": 12, "color": DEEP}),
    ], line_spacing=1.3)
    # IT right
    rounded_box(slide, 6.83, col_y, col_w, col_h, fill=SURFACE, stroke=TEAL, stroke_w=2.0)
    rectangle(slide, 6.83, col_y, col_w, 0.6, fill=TEAL)
    text_box(slide, 6.83, col_y + 0.05, col_w, 0.5, "IT · Information Technology",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 7.08, col_y + 0.8, col_w - 0.5, 3.5, [
        ("Где живёт:", {"size": 11, "bold": True, "color": LIGHT}),
        ("Облако · дата-центры · корпоративные базы · ML pipelines", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Цикл:", {"size": 11, "bold": True, "color": LIGHT}),
        ("100–500 мс недетерминированный (LLM)", {"size": 13, "bold": True, "color": GOLD}),
        ("", {"size": 6}),
        ("Консистентность: eventually-consistent", {"size": 12, "color": DEEP}),
        ("Audit trail: security и compliance, не control", {"size": 12, "color": DEEP}),
        ("Sertification: SaaS-обновления каждую неделю", {"size": 12, "color": DEEP}),
    ], line_spacing=1.3)
    # Conclusion
    rounded_box(slide, 0.5, 6.4, 12.33, 0.7, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 0.7, 6.45, 12.0, 0.6,
             "LLM с задержкой 100–500 мс физически не вмещается в PLC-цикл 1–10 мс — это арифметика, не временное препятствие",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "OT — мир PLC, SCADA, конвейеров. Время в миллисекундах, детерминированный budget. Audit trail регуляторный. SIL 2/3 для safety. IT — облако, ML pipelines. Время в миллисекундах-секундах, eventually-consistent. AI родился в IT-мире и приходит в OT. Упирается в structural divide. LLM 100-500 мс не вмещается в PLC-цикл 1-10 мс. Это арифметика. Из этого: AI на L3-L4 ISA-95 работает. На L2 — со скрипом через decision-support. На L1 — почти никогда. На L0 — никогда.")


def s10_foundation_models(p):
    """s10 — foundation models = augmentation, 3 reasons."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Фундаментальные модели = augmentation, НЕ controller",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Siemens IFM + Foxconn FoxBrain — три причины, почему не контроллер",
             size=14, italic=True, color=LIGHT)
    # Left: examples
    rounded_box(slide, 0.5, 1.7, 4.0, 5.1)
    text_box(slide, 0.7, 1.85, 3.6, 0.4, "ЧТО ЕСТЬ",
             size=14, bold=True, color=MID)
    multiline_box(slide, 0.7, 2.3, 3.6, 2.4, [
        ("Siemens IFM", {"size": 14, "bold": True, "color": DEEP}),
        ("Hannover Messe март 2025", {"size": 10, "italic": True, "color": SLATE}),
        ("150 ПБ engineering data", {"size": 11, "color": DEEP}),
        ("1000+ AI-патентов Siemens", {"size": 11, "color": DEEP}),
        ("", {"size": 10}),
        ("Foxconn FoxBrain", {"size": 14, "bold": True, "color": DEEP}),
        ("Computex 2025", {"size": 10, "italic": True, "color": SLATE}),
        ("Llama 3.1 70B derivative", {"size": 11, "color": DEEP}),
        ("Injection-molding параметры", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    # Siemens HQ image bottom
    img_path = ASSETS / "screenshots" / "s10-siemens-hq.jpg"
    if img_path.exists():
        add_image(slide, img_path, 0.7, 4.85, 3.6, 1.7)
        attribution(slide, "Siemens HQ Munich · Wikimedia", x=0.7, y=6.55, w=3.6)
    # Right: 3 reasons
    multiline_box(slide, 4.8, 1.85, 8.1, 0.5, [
        ("ЧТО НЕ ДЕЛАЮТ — три причины", {"size": 14, "bold": True, "color": MID}),
    ])
    reasons = [
        ("1. Задержка вывода", "LLM 100–500 мс vs PLC-цикл 1–10 мс. Физически не вмещается в control loop.", MID),
        ("2. Галлюцинации", "Недетерминированный output несовместим с замыканием контура управления.", LIGHT),
        ("3. Сертификация", "Нет audit trail для SIL 2/3, FDA Part 11, GAMP®5. Black-box ML не проходит валидацию.", TEAL),
    ]
    for i, (title, desc, color) in enumerate(reasons):
        y = 2.4 + i * 1.35
        rounded_box(slide, 4.8, y, 8.1, 1.2)
        rectangle(slide, 4.8, y, 0.15, 1.2, fill=color)
        text_box(slide, 5.05, y + 0.15, 7.85, 0.45, title,
                 size=15, bold=True, color=DEEP)
        text_box(slide, 5.05, y + 0.65, 7.85, 0.55, desc,
                 size=12, color=DARK_GREY, line_spacing=1.3)
    # Conclusion
    rounded_box(slide, 4.8, 6.5, 8.1, 0.6, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 5.0, 6.55, 7.9, 0.5,
             "Chat-помощник для оператора ≠ autonomous controller — два разных архитектурных класса",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "Siemens IFM на Hannover Messe март 2025 — обучен на 150 ПБ engineering data + 1000 AI-патентов. Foxconn FoxBrain март 2025 — derivative Llama 3.1 70B с injection-molding параметрами. Это серьёзные модели от серьёзных компаний. И вокруг маркетинг про революцию. Я хочу, чтобы вы сделали один концептуальный шаг — это augmentation, не autonomous controller. Три причины. Задержка вывода — 100-500 мс не вмещается в PLC. Галлюцинации — недетерминированный output несовместим с control loop. Сертификация — нет audit trail для SIL 2/3, FDA Part 11. Эти причины не уходят со временем.")


def s11_optimus_reality(p):
    """s11 — Tesla Optimus reality check."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Tesla Optimus: демо есть, production нет",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Recurring pattern: capability в контролируемой среде vs reliability 99,9% в неконтролируемой",
             size=14, italic=True, color=LIGHT)
    # Two columns
    col_y = 1.8
    col_h = 4.5
    col_w = 6.0
    # DEMO left
    rounded_box(slide, 0.5, col_y, col_w, col_h, fill=SURFACE, stroke=LIGHT, stroke_w=2.0)
    rectangle(slide, 0.5, col_y, col_w, 0.6, fill=LIGHT)
    text_box(slide, 0.5, col_y + 0.05, col_w, 0.5, "ДЕМО · 2022–2024",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.75, col_y + 0.8, col_w - 0.5, 3.5, [
        ("Tesla AI Day 2022, 2024 · Cybercab event Oct 2024", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("Видеодемо:", {"size": 13, "bold": True, "color": MID}),
        ("· Складывает футболки", {"size": 13, "color": DEEP}),
        ("· Переносит коробки", {"size": 13, "color": DEEP}),
        ("· Ходит по фабрике", {"size": 13, "color": DEEP}),
        ("· Передаёт инструменты", {"size": 13, "color": DEEP}),
        ("", {"size": 8}),
        ("Маркетинг:", {"size": 13, "bold": True, "color": MID}),
        ("10–20 тыс. Optimus к 2025, цена ~$30K, миллионы к 2030", {"size": 12, "color": DEEP}),
    ], line_spacing=1.3)
    # REALITY right
    rounded_box(slide, 6.83, col_y, col_w, col_h, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    rectangle(slide, 6.83, col_y, col_w, 0.6, fill=GOLD)
    text_box(slide, 6.83, col_y + 0.05, col_w, 0.5, "PRODUCTION · 2026",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 7.08, col_y + 0.8, col_w - 0.5, 3.5, [
        ("На февраль 2026:", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("Tesla Optimus НЕ задействован в production-сборке ни на одном заводе.", {"size": 14, "bold": True, "color": DEEP}),
        ("", {"size": 8}),
        ("Cybercab event Oct 2024:", {"size": 12, "bold": True, "color": MID}),
        ("Optimus раздавал напитки гостям — Wired подтвердил, что роботы управлялись людьми удалённо.", {"size": 12, "italic": True, "color": DEEP}),
    ], line_spacing=1.35)
    # Conclusion
    rounded_box(slide, 0.5, 6.45, 12.33, 0.55, fill=SURFACE, stroke=LIGHT)
    text_box(slide, 0.7, 6.5, 12.0, 0.45,
             "Между демо и production — обычно 5–10 лет, и большинство демо не переходят",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "Tesla Optimus — отличный пример паттерна демо-vs-production. Маркетинг яркий: 10-20 тысяч Optimus к 2025, миллионы к 2030. Reality check февраль 2026: не задействован в production-сборке ни на одном заводе. На Cybercab event октябрь 2024 Optimus раздавал напитки — Wired подтвердил, роботы управлялись людьми удалённо. Это не критика Tesla — это иллюстрация паттерна. Демо демонстрирует capability в контролируемой среде. Production — надёжность 99,9% в неконтролируемой × десятки тысяч часов × тысячи юнитов. Между этими состояниями обычно 5-10 лет.")


def s12_hype_collapse_trio(p):
    """s12 — three hype-collapses."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Три истории на $4B+ каждая. Один урок: AI ≠ магия.",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "GE Predix · IBM Watson · Foxconn Wisconsin — три провала, три урока",
             size=14, italic=True, color=LIGHT)
    cards = [
        ("GE Predix", "2011–2020", "$4B+",
         "AOS for industrial — свой облачный конкурент AWS.\n«be everything to everyone» → развал 2018, продажа частями к 2020.",
         "Industrial AI ≠ general cloud AI. Размер инвестиций ≠ результат.", MID),
        ("IBM Watson", "2018–2022", "Multi-$B",
         "«–47% downtime, –48% defect rate» (маркетинг).\nWatson for Drug Discovery остановлен 2019. Watson Health продан за $1B в 2022.",
         "Демо ≠ production. Маркетинговые проценты ≠ measured production metrics.", TEAL),
        ("Foxconn Wisconsin", "2018–2024", "$10B обещано",
         "«Восьмое чудо света» (2018).\n10 000 рабочих → 1 500. LCD → screens → AI hub → coffee kiosks → Microsoft Fairwater data center ($3,3B, май 2024).",
         "«Чудо света» от главы государства — anti-signal. Маркетинг ≠ физика.", GOLD),
    ]
    card_w = 4.05
    gap = 0.1
    for i, (name, dates, money, body, lesson, color) in enumerate(cards):
        x = 0.5 + i * (card_w + gap)
        y = 1.7
        rounded_box(slide, x, y, card_w, 5.2)
        # Header strip
        rectangle(slide, x, y, card_w, 0.7, fill=color)
        text_box(slide, x + 0.1, y + 0.05, card_w - 0.2, 0.6, name,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Dates + money
        text_box(slide, x + 0.15, y + 0.85, card_w - 0.3, 0.3, dates,
                 size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
        text_box(slide, x + 0.15, y + 1.2, card_w - 0.3, 0.5, money,
                 size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        # Body
        text_box(slide, x + 0.25, y + 1.85, card_w - 0.5, 2.2, body,
                 size=11, color=DEEP, line_spacing=1.3)
        # Lesson
        rounded_box(slide, x + 0.2, y + 4.2, card_w - 0.4, 0.9,
                    fill=GOLD_TINT, stroke=GOLD, stroke_w=1.0)
        text_box(slide, x + 0.3, y + 4.25, card_w - 0.6, 0.8, "Урок: " + lesson,
                 size=11, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    attribution(slide, "GE 10-K filings 2017-2020 · IBM Watson Health 2022 sale press · Microsoft Fairwater press May 2024")
    add_notes(slide, "Три истории, каждая 4 миллиарда плюс. GE Predix: с 2011 по 2020 GE потратил 4 миллиарда плюс на промышленное облако, конкурент AWS. «Be everything to everyone» — провал. Урок: industrial AI ≠ general cloud AI. IBM Watson: маркетинг 47 процентов снижение downtime, Watson for Drug Discovery остановлен 2019, Watson Health продан за миллиард в 2022. Урок: демо ≠ production. Foxconn Wisconsin: «восьмое чудо света» в 2018, 10 тысяч рабочих обещано → 1500 к 2024, площадка стала дата-центром Microsoft. Урок: «чудо света» от главы государства — anti-signal.")


def apply_md_notes(pres):
    """Read slides/*.md and replace generated notes with full chapter-derived notes."""
    slides_dir = ROOT.parent / "slides"
    if not slides_dir.exists():
        return
    md_files = sorted(slides_dir.glob("s*.md"))
    for i, slide in enumerate(pres.slides):
        if i >= len(md_files):
            break
        md = md_files[i].read_text(encoding="utf-8")
        m = re.search(r"## Speaker notes\s*\n(.*?)(?:\n##|\Z)", md, re.S)
        if m:
            notes = m.group(1).strip()
            paragraphs = [p.strip().replace("\n", " ") for p in notes.split("\n\n") if p.strip()]
            notes_text = "\n\n".join(paragraphs)
            try:
                slide.notes_slide.notes_text_frame.text = notes_text
            except Exception as e:
                print(f"Notes apply fail slide {i+1}: {e}")


# Import part 2 for the remaining slides
import sys
sys.path.insert(0, str(ROOT))


if __name__ == "__main__":
    from build_lec11_part2 import (
        s13_section2_divider, s14_cv_inspection, s15_boeing_737,
        s16_label_cost, s17_pdm_oee, s18_cobots_jidoka,
        s19_tesla_2018, s20_cv_limits, s21_foxbrain, s22_discrete_matrix,
        s23_section3_divider, s24_soft_sensors, s25_mpc_rl_cirl,
        s26_rl_drift, s27_edge_pdm, s28_regulatory, s29_russian, s30_process_matrix,
        s31_section4_divider, s32_criteria, s33_alternatives, s34_pfizer_vox, s35_framework,
        s36_section5_divider, s37_recap, s38_qa, s39_closing,
    )

    pres = setup_pres()
    # Section 0
    s01_hero_tesla(pres)
    s02_cover(pres)
    s03_lecture_map(pres)
    s04_glossary(pres)
    s05_keystone(pres)
    # Section 1
    s06_section1_divider(pres)
    s07_adoption(pres)
    s08_market_estimates(pres)
    s09_ot_it_split(pres)
    s10_foundation_models(pres)
    s11_optimus_reality(pres)
    s12_hype_collapse_trio(pres)
    # Section 2
    s13_section2_divider(pres)
    s14_cv_inspection(pres)
    s15_boeing_737(pres)
    s16_label_cost(pres)
    s17_pdm_oee(pres)
    s18_cobots_jidoka(pres)
    s19_tesla_2018(pres)
    s20_cv_limits(pres)
    s21_foxbrain(pres)
    s22_discrete_matrix(pres)
    # Section 3
    s23_section3_divider(pres)
    s24_soft_sensors(pres)
    s25_mpc_rl_cirl(pres)
    s26_rl_drift(pres)
    s27_edge_pdm(pres)
    s28_regulatory(pres)
    s29_russian(pres)
    s30_process_matrix(pres)
    # Section 4
    s31_section4_divider(pres)
    s32_criteria(pres)
    s33_alternatives(pres)
    s34_pfizer_vox(pres)
    s35_framework(pres)
    # Section 5
    s36_section5_divider(pres)
    s37_recap(pres)
    s38_qa(pres)
    s39_closing(pres)

    print(f"Slides built: {len(pres.slides)}")

    apply_md_notes(pres)

    pres.save(str(OUT))
    print(f"Saved: {OUT}")
