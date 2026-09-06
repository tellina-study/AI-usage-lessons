"""
Reference / page-number system for Лекции 3 (issue #171, ported from
lec-04 `_helpers.py` ПРАВКА 1 + #170-3/#170-4 techniques).

Provides:
  • URLS            — canonical URL map (ONLY from reference-registry.md).
  • SLIDE_REFS      — per display-slide source registry:
                        (num, short_name, urlkey|None, gloss[, volatile]).
  • ANCHORS         — per display-slide list of (ref_nums, anchor_substr):
                        where to inject the small superscript [N] marker(s)
                        inside the already-built visible body. anchor_substr
                        is a verbatim fragment of an EXISTING run — nothing in
                        the visible copy changes except the appended [N].
  • shrink_refs_in_frame  — post-hoc split of [N] markers into small (~52%)
                            superscript muted runs (#170-3).
  • inject_ref_markers    — walk a slide, append [N] at each ANCHORS anchor,
                            then shrink.
  • ref_list / refs_of_slide  — bottom clickable numbered source list.
  • notes_sources_block / notes_with_sources — «Источники:» block for notes,
                            [VFY]/[VFY-day-of] on volatile/unconfirmed sources.
  • page_number     — muted «N / 40» stamp bottom-right.

[VFY] policy (task req 4): course-internal / illustrative / future-dated
sources (agent-harness-registry; arXiv:2605.29420 / 2605.29463 / 2603.22489 /
2601.06007 / 2601.18699) are NEVER given as canonical primary URLs on the
visible slide — the confirmed anchor (Zheng, Gloaguen, Willison/Docker/Unit42,
Luo …) carries the [N]. The unconfirmed ones are named ONLY in notes with a
[VFY] tag. Volatile-but-real numbers get [VFY-day-of] in notes.
"""
import re

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# palette (mirror build_v3.py)
DEEP  = RGBColor(0x21, 0x29, 0x5C)
MID   = RGBColor(0x06, 0x5A, 0x82)
LIGHT = RGBColor(0x1C, 0x72, 0x93)
SLATE = RGBColor(0x5B, 0x66, 0x78)
FONT_BODY = "Arial"

# ============================================================
# URLS — canonical, ONLY from notes/research/lecture-3/reference-registry.md
# ============================================================
URLS = {
    # Air Canada / legal
    "aircanada_mccarthy": "https://www.mccarthy.ca/en/insights/blogs/techlex/moffatt-v-air-canada-misrepresentation-ai-chatbot",
    "aircanada_aba": "https://www.americanbar.org/groups/business_law/resources/business-law-today/2024-february/bc-tribunal-confirms-companies-remain-liable-information-provided-ai-chatbot/",
    # Anthropic engineering / research
    "anthropic_agents": "https://www.anthropic.com/research/building-effective-agents",
    "anthropic_context_eng": "https://www.anthropic.com/engineering/effective-context-engineering-for-ai-agents",
    "anthropic_cot_faith": "https://www.anthropic.com/research/reasoning-models-dont-say-think",
    "anthropic_retention": "https://platform.claude.com/docs/en/build-with-claude/api-and-data-retention",
    "anthropic_mcp_donate": "https://www.anthropic.com/news/donating-the-model-context-protocol-and-establishing-of-the-agentic-ai-foundation",
    # Prompt / roles / CoT
    "zheng_personas": "https://arxiv.org/abs/2311.10054",
    "wei_cot": "https://arxiv.org/abs/2201.11903",
    # Context
    "chroma_context_rot": "https://research.trychroma.com/context-rot",
    "liu_lost_middle": "https://arxiv.org/abs/2307.03172",
    # RAG
    "lewis_rag": "https://arxiv.org/abs/2005.11401",
    "barnett_7fail": "https://arxiv.org/abs/2401.05856",
    "kore_7rag": "https://www.kore.ai/blog/seven-rag-engineering-failure-points",
    "uniah": "https://arxiv.org/abs/2503.00353",
    "redhat_rag_ft": "https://www.redhat.com/en/topics/ai/rag-vs-fine-tuning",
    "ibm_rag_ft": "https://www.ibm.com/think/topics/rag-vs-fine-tuning",
    "bigdata_ft": "https://bigdataboutique.com/blog/fine-tuning-llms-when-rag-isnt-enough",
    # Fine-tuning / PEFT / distillation
    "hf_beyond_lora": "https://huggingface.co/blog/peft-beyond-lora",
    "lora_paper": "https://arxiv.org/abs/2106.09685",
    "qlora_paper": "https://arxiv.org/abs/2305.14314",
    "hinton_distill": "https://arxiv.org/abs/1503.02531",
    "peft_survey": "https://arxiv.org/abs/2403.14608",
    "luo_forgetting": "https://arxiv.org/abs/2308.08747",
    # Agents / loop
    "yao_react": "https://arxiv.org/abs/2210.03629",
    "cognition_no_multiagent": "https://cognition.ai/blog/dont-build-multi-agents",
    "gloaguen_presence": "https://arxiv.org/abs/2602.11988",
    "claude_code_51735": "https://github.com/anthropics/claude-code/issues/51735",
    # Security / retention
    "willison_mcp_inject": "https://simonwillison.net/2025/Apr/9/mcp-prompt-injection/",
    "docker_mcp_horror": "https://www.docker.com/blog/mcp-horror-stories-github-prompt-injection/",
    "unit42_mcp": "https://unit42.paloaltonetworks.com/model-context-protocol-attack-vectors/",
    "nyt_openai_bloomberg": "https://news.bloomberglaw.com/ip-law/openai-must-turn-over-20-million-chatgpt-logs-judge-affirms",
    # Failures / postmortems / market
    "jain_4200": "https://medium.com/@sattyamjain96/the-agent-that-burned-4-200-in-63-hours-a-production-ai-postmortem-d38fd9586a85",
    "mindstudio_reliability": "https://www.mindstudio.ai/blog/reliability-compounding-problem-ai-agent-stacks",
    "mit_nanda_fortune": "https://fortune.com/2025/08/18/mit-report-95-percent-generative-ai-pilots-at-companies-failing-cfo/",
}

# ============================================================
# SLIDE_REFS — (num, short_name, urlkey|None, gloss[, volatile])
# volatile=True → [VFY-day-of] in notes. urlkey=None → course-internal
# (agent-harness-registry, unconfirmed) → notes-only [VFY], never on slide.
# ============================================================
SLIDE_REFS = {
    "s01b": [
        ("1", "McCarthy Tétrault — Moffatt v. Air Canada (BC CRT, 14.02.2024)", "aircanada_mccarthy",
         "бот выдумал политику возврата; трибунал: компания отвечает за ответ бота"),
        ("2", "ABA Business Law Today — компании отвечают за AI-чат-бота", "aircanada_aba",
         "«бот — не отдельное юр. лицо»: неправильный выбор архитектуры под задачу lookup"),
    ],
    "s05": [
        ("1", "Anthropic — Building Effective Agents («найди простейшее»)", "anthropic_agents",
         "не усложняй архитектуру без требования задачи — распределение бремени доказательства"),
    ],
    "s05a": [
        ("1", "Zheng et al. 2024, Findings of EMNLP — персоны не повышают точность фактов", "zheng_personas",
         "162 персоны · 8 доменов · 2410 фактических вопросов · 4 семейства → нет прироста точности"),
    ],
    "s06": [
        ("1", "Wei et al. 2022 — Chain-of-Thought Prompting", "wei_cot",
         "пошаговое рассуждение поднимает надёжность на арифметике/многошаговой логике"),
        ("2", "Anthropic — Reasoning Models Don't Always Say What They Think", "anthropic_cot_faith",
         "faithfulness: Claude 3.7 ~25%, DeepSeek R1 ~39% — рассуждение не обязано отражать причину", True),
    ],
    "s08": [
        ("1", "Chroma Research — Context Rot", "chroma_context_rot",
         "точность извлечения падает с ростом числа токенов в контексте"),
        ("2", "Liu et al. 2023 — Lost in the Middle", "liu_lost_middle",
         "тот же феномен «lost in the middle» из Лекции 2 — новый термин, не новая сущность"),
        ("3", "Anthropic — Effective Context Engineering", "anthropic_context_eng",
         "минимальный высокосигнальный контекст — инженерное требование, не эстетика"),
    ],
    "s10": [
        ("1", "Lewis et al. 2020 — Retrieval-Augmented Generation (NeurIPS)", "lewis_rag",
         "каноническая статья RAG: индексация → retrieval → генерация с опорой"),
    ],
    "s11": [
        ("1", "IBM — RAG vs Fine-Tuning (вендор-нейтрально)", "ibm_rag_ft",
         "признаки-за-RAG: большое/меняется/провенанс/приватное"),
        ("2", "U-NIAH — RAG win-rate выше у меньших моделей", "uniah",
         "выигрыш RAG над прямым ответом особенно велик для меньших моделей", True),
    ],
    "s12": [
        ("1", "Red Hat — RAG vs Fine-Tuning (когда не RAG)", "redhat_rag_ft",
         "корпус влезает в окно → full-context+кэш; фикс. значение → lookup; live → API без индекса"),
        ("2", "McCarthy Tétrault — Air Canada («генерация поверх фикс. политики»)", "aircanada_mccarthy",
         "фиксированная политика → детерминированный lookup, не retrieval+генерация"),
    ],
    "s13": [
        ("1", "Barnett et al. 2024 — Seven Failure Points (RAG)", "barnett_7fail",
         "«вернул что-то ≠ вернул правильное»: 7 точек отказа RAG-инжиниринга"),
        ("2", "Kore.ai — Seven RAG Engineering Failure Points", "kore_7rag",
         "legal-AI / medical-RAG / support-бот — деградация на масштабе без observability"),
        ("3", "McCarthy Tétrault — Air Canada (отказ grounding)", "aircanada_mccarthy",
         "сгенерированный текст в роли, требовавшей извлечённого проверенного факта"),
    ],
    "s13b": [
        ("1", "IBM — RAG vs Fine-Tuning (fine-tuning меняет веса)", "ibm_rag_ft",
         "промпт/RAG меняют контекст; fine-tuning меняет сами веса модели"),
    ],
    "s14": [
        ("1", "Hinton, Vinyals, Dean 2015 — Distilling the Knowledge", "hinton_distill",
         "дистилляция — самостоятельная техника, таксономически НЕ вид fine-tuning"),
        ("2", "PEFT survey — таксономия методов настройки", "peft_survey",
         "обзоры разносят дистилляцию и fine-tuning по разным категориям"),
    ],
    "s17": [
        ("1", "BigData Boutique — Fine-Tuning When RAG Isn't Enough", "bigdata_ft",
         "fine-tuning сузился до поведения/стиля/формата/политики; знание → RAG", True),
    ],
    "s15": [
        ("1", "Hu et al. 2021 — LoRA", "lora_paper",
         "базовые веса замораживаются, обучаются низкоранговые адаптеры"),
        ("2", "Dettmers et al. 2023 — QLoRA", "qlora_paper",
         "LoRA поверх квантованной модели → дообучение на одном GPU"),
        ("3", "HF PEFT — Beyond LoRA (LoRA 98,4% из 20 834 карточек)", "hf_beyond_lora",
         "98,4% моделей с тегом PEFT используют LoRA; оговорка: доля среди помеченных PEFT", True),
    ],
    "s16": [
        ("1", "Luo et al. 2023 — Catastrophic Forgetting in LLM Continual FT", "luo_forgetting",
         "узкий агрессивный FT ломает общие способности; тяжелее с ростом масштаба модели"),
    ],
    "s19": [
        ("1", "Anthropic — MCP donation / Agentic AI Foundation (N×M→N+M)", "anthropic_mcp_donate",
         "MCP стандартизует подключение; удобство подключения ≠ безопасность подключаемого", True),
    ],
    "s21": [
        ("1", "Yao et al. 2022 — ReAct (plan→act→check→iterate)", "yao_react",
         "чередование рассуждения и действий; каждый шаг — место отказа"),
        ("2", "Anthropic — Reasoning Models faithfulness (check ≠ самооценка)", "anthropic_cot_faith",
         "шаг check — валидация против внешнего критерия, не самооценка модели", True),
    ],
    "s22": [
        ("1", "Anthropic — Building Effective Agents (workflow vs agent)", "anthropic_agents",
         "workflow = предопределённые пути; агент = динамический процесс; размен латентность/стоимость↔качество"),
        ("2", "Cognition — Don't Build Multi-Agents", "cognition_no_multiagent",
         "мульти-агент по умолчанию не апгрейд; хрупкость параллельных субагентов"),
    ],
    "s22b": [
        ("1", "agent-harness-registry — карта слотов экипировки агента", None,
         "источник не подтверждён canonical-URL 2026-08-30; Claude Code/Cursor/Aider — вендор-сайты, verify day-of", True),
    ],
    "s22c": [
        ("1", "agent-harness-registry (live-eval) — спектр памяти агента", None,
         "источник не подтверждён canonical-URL 2026-08-30; параллель с критерием масштаба RAG", True),
    ],
    "s22d": [
        ("1", "agent-harness-registry (live-eval) — Letta Tier D / Memory Tool 17%-хвост", None,
         "источник не подтверждён canonical-URL 2026-08-30; числа volatile (1.0/0.833/0.750; 17%)", True),
    ],
    "s22e": [
        ("1", "Gloaguen et al. 2026 — Evaluating AGENTS.md (presence paradox)", "gloaguen_presence",
         "наличие файла-инструкции не даёт значимого прироста; помогает в пробеле документации", True),
        ("2", "GitHub anthropics/claude-code#51735 (повтор ошибки через 25 дней)", "claude_code_51735",
         "письменная запись о прошлой ошибке не предотвратила её повтор"),
    ],
    "s23": [
        ("1", "Sattyam Jain 2026 — The Agent That Burned $4,200 in 63 Hours", "jain_4200",
         "петля без лимитов на HTTP 429; retry-скрипт решил бы задачу почти бесплатно", True),
        ("2", "MindStudio — Reliability Compounding Problem", "mindstudio_reliability",
         "5×99%≈95%, 10→90%, 20→82% — надёжности перемножаются"),
        ("3", "Cognition — Don't Build Multi-Agents (хрупкость)", "cognition_no_multiagent",
         "зависимые подзадачи → параллельные субагенты принимают конфликтующие решения"),
    ],
    "s25": [
        ("1", "Docker — MCP Horror Stories: GitHub Prompt Injection", "docker_mcp_horror",
         "GitHub MCP heist: issue-инструкция + широкий токен → выгрузка приватных репо"),
        ("2", "Simon Willison — Prompt injection via MCP", "willison_mcp_inject",
         "модель не отличает данные от команды; недоверенный контент = команда"),
        ("3", "Palo Alto Unit 42 — MCP Attack Vectors", "unit42_mcp",
         "tool poisoning / каждое подключение = новая граница доверия"),
        ("4", "Bloomberg Law — NYT v. OpenAI (суд обязал хранить логи)", "nyt_openai_bloomberg",
         "ZDR покрывает не всё; судебный приказ поверх любой политики хранения"),
        ("5", "Anthropic — API and Data Retention (границы ZDR)", "anthropic_retention",
         "ZDR не покрывает third-party / MCP-connector — ровно то, из чего агент состоит", True),
    ],
    "s25b": [
        ("1", "agent-harness-registry — обзор через рамку экипировки", None,
         "источник не подтверждён; Claude Code/Aider/Cursor/OpenHands — вендор-сайты, verify day-of", True),
    ],
    "s26": [
        ("1", "Anthropic — Building Effective Agents (правило лестницы)", "anthropic_agents",
         "оставайся на нижней ступени; каждый подъём — обмен, не улучшение"),
    ],
    "s27": [
        ("1", "Anthropic — Building Effective Agents (маршрут выбора)", "anthropic_agents",
         "маршрут вопросов сверху вниз; шаг 1 — детерминированная задача → обычный код, СТОП"),
        ("2", "McCarthy Tétrault — Air Canada (нижняя строка матрицы)", "aircanada_mccarthy",
         "генеративная архитектура на детерминированную задачу = ошибка нижней строки"),
    ],
    "s29": [
        ("1", "Anthropic — Reasoning Models faithfulness (self-rationale ≠ контроль)", "anthropic_cot_faith",
         "человек проверяет результат/факты против источника, не самообъяснение модели", True),
        ("2", "MIT NANDA — State of AI in Business 2025 (~95% пилотов без ROI)", "mit_nanda_fortune",
         "корень — learning gap и провал интеграции, не качество модели; отчёт, не закон", True),
    ],
}

# ============================================================
# ANCHORS — where [N] markers go inside the already-built visible body.
# (ref_nums:str, anchor_substr:str). anchor_substr MUST be a verbatim
# fragment of an existing run; the marker «[N]» is appended right after it.
# Nothing else in the visible copy changes.
# NB: s22c/s22d/s25b anchors carry a course-internal source (urlkey None) →
# still get a visible [N] whose bottom-list entry has NO hyperlink (framed as
# «независимый реестр live-eval», with [VFY] only in notes).
# ============================================================
ANCHORS = {
    "s01b": [("1,2", "Трибунал: «бот — не отдельное юр. лицо» → компания вернула $812,02")],
    "s05": [("1", "Не усложняй архитектуру без причины, выраженной в требованиях задачи")],
    "s05a": [("1", "персоны НЕ повысили точность")],
    "s06": [("1", "23 − 7 = 16"),
            ("2", "не обязано отражать реальную причину ответа")],
    "s08": [("1,2", "context rot = тот же «lost in the middle» из Л2"),
            ("3", "это инженерное требование, не эстетика")],
    "s10": [("1", "«Не знаю» / «см. источник X» — корректный ответ RAG-системы")],
    "s11": [("1", "RAG оправдан при сильном сигнале по признакам ниже"),
            ("2", "Один признак — повод присмотреться")],
    "s12": [("1", "RAG избыточен, если выполнен ЛЮБОЙ из трёх"),
            ("2", "Фиксированная политика / значение")],
    "s13": [("1,2", "«Вернул что-то» ≠ «вернул правильное»"),
            ("3", "Air Canada — разбор архитектуры")],
    "s13b": [("1", "Дообучение = «изменить саму модель»")],
    "s14": [("1,2", "две таксономически разные операции")],
    "s17": [("1", "гибрид — норма")],
    "s15": [("1,2", "LoRA — низкоранговые матрицы-адаптеры; QLoRA — то же поверх квантованной модели"),
            ("3", "моделей с тегом PEFT — это LoRA")],
    "s16": [("1", "деградация общих способностей модели в результате узкого агрессивного дообучения")],
    "s19": [("1", "N+M")],
    "s21": [("1", "Агент — архитектура, где модель не делает один проход, а работает в цикле"),
            ("2", "валидация против ВНЕШНЕГО критерия — не самооценка модели")],
    "s22": [("1", "Предсказуемая задача → сценарий"),
            ("2", "ценность оправдывает кратные стоимость/риск → агент")],
    "s22b": [("1", "Пять типовых слотов")],
    "s22c": [("1", "Спектр — от плоского файла до граф-баз")],
    "s22d": [("1", "Независимая проверка показывает: иногда — драматически нет")],
    "s22e": [("1", "RCT (Gloaguen et al. 2026): само наличие файла-инструкции НЕ даёт значимого прироста"),
             ("2", "письменно признанная прошлая ошибка НЕ предотвратила её повторение спустя 25 дней")],
    "s23": [("1", "$4 200 за 63 часа"),
            ("2", "«меньше переходов + проверка между шагами» — сильный"),
            ("3", "зависимые подзадачи → параллельные субагенты принимают конфликтующие")],
    "s25": [("1", "ассистент выгрузил приватные репозитории в публичный PR"),
            ("2,3", "как только агент делегирует и подключается — появляется поверхность"),
            ("4,5", "судебный приказ (NYT v. OpenAI) + сторонние сервисы/MCP вне ZDR")],
    "s25b": [("1", "какие слоты экипировки заполнены и где агент физически живёт")],
    "s26": [("1", "Каждый подъём — это ОБМЕН")],
    "s27": [("1", "останавливайтесь на первом сработавшем вопросе"),
            ("2", "если задача детерминированная и верифицируемая — обычный код, СТОП здесь")],
    "s29": [("1", "не по правдоподобности рассуждения"),
            ("2", "корень в разрыве обучения и провале интеграции, не в качестве модели")],
}

# ============================================================
# NOTES_ANCHORS — where [N] markers go inside the READABLE speaker notes
# (.md «## Speaker notes»). (ref_nums:str, notes_anchor_substr:str) — the
# marker «[N]» is appended right after the first verbatim occurrence.
# ============================================================
NOTES_ANCHORS = {
    "s01b": [("1,2", "отвечающее за свои действия».")],
    "s05": [("1", "один вызов модели с хорошо составленным промптом")],
    "s05a": [("1", "сдвигает ли это фактическую точность")],
    "s06": [("1", "16 плюс 12 равно 28"),
            ("2", "примерно в двух из пяти")],
    "s08": [("3", "видимых модели на инференсе")],
    "s10": [("1", "опираясь на эти фрагменты")],
    "s11": [("1", "сильный сигнал по одному или нескольким признакам")],
    "s12": [("1", "Три явных критерия «не RAG»")],
    "s13": [("1,2", "не означает «система вернула правильное»"),
            ("3", "Air Canada")],
    "s13b": [("1", "продолжение обучения уже готовой, предобученной модели")],
    "s14": [("1,2", "дистилляц")],
    "s17": [("1", "гибрид")],
    "s15": [("1", "адаптеров"),
            ("3", "98")],
    "s16": [("1", "деградация общих способностей модели")],
    "s19": [("1", "открытый стандарт единого способа подключать")],
    "s21": [("1", "работает в цикле"),
            ("2", "шаг check")],
    "s22": [("1", "разделяет два понятия"),
            ("2", "мульти-агент")],
    "s22b": [("1", "цикл плюс оснастка")],
    "s22c": [("1", "плоский файл")],
    "s22d": [("1", "независимого реестра")],
    "s22e": [("1", "presence paradox"),
             ("2", "claude-code")],
    "s23": [("1", "429"),
            ("2", "перемножа"),
            ("3", "мульти-агент")],
    "s25": [("1", "приватные репозитории и опубликуй здесь»"),
            ("2,3", "не отличает данные от команды"),
            ("4,5", "судебный приказ")],
    "s25b": [("1", "рамку экипировки")],
    "s26": [("1", "лестницу архитектурной сложности")],
    "s27": [("1", "останавливаясь на первом сработавшем"),
            ("2", "детерминированная и верифицируемая")],
    "s29": [("1", "unfaithful"),
            ("2", "learning gap")],
}


def inject_notes_markers(note_text, sid):
    """Insert [N] markers into readable notes at NOTES_ANCHORS anchors.
    Returns (new_text, missed:list[str])."""
    missed = []
    for ref_nums, anchor in NOTES_ANCHORS.get(sid, []):
        marker = f"[{ref_nums}]"
        if marker in note_text:
            continue
        i = note_text.find(anchor)
        if i < 0:
            missed.append(anchor)
            continue
        j = i + len(anchor)
        note_text = note_text[:j] + marker + note_text[j:]
    return note_text, missed


# ============================================================
# [N] shrink (#170-3): split [N] markers into small superscript muted runs.
# ============================================================
_REF_RE = re.compile(r'\[\d+(?:\s*[,–—-]\s*\d+)*\]')
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _run_props(src_run):
    f = src_run.font
    sz = f.size
    return {
        "name": f.name,
        "size_pt": (sz.pt if sz is not None else None),
        "bold": f.bold,
        "italic": f.italic,
        "color": (f.color.rgb if (f.color and f.color.type is not None) else None),
    }


def _clone_run_after(anchor_r, props, text, *, ref=False, ref_frac=0.52,
                     ref_color=LIGHT):
    new_r = etree.SubElement(anchor_r.getparent(), _A + "r")
    anchor_r.addnext(new_r)
    rpr = etree.SubElement(new_r, _A + "rPr")
    base = props["size_pt"] or 16.0
    if ref:
        rpr.set("sz", str(int(round(base * ref_frac * 100))))
        rpr.set("baseline", "30000")
        rpr.set("b", "0")
        rpr.set("i", "1")
    else:
        if props["size_pt"] is not None:
            rpr.set("sz", str(int(round(base * 100))))
        if props["bold"] is not None:
            rpr.set("b", "1" if props["bold"] else "0")
        if props["italic"] is not None:
            rpr.set("i", "1" if props["italic"] else "0")
    if props["name"]:
        for tag in ("latin", "cs", "ea"):
            el = etree.SubElement(rpr, _A + tag)
            el.set("typeface", props["name"])
    col = ref_color if ref else props["color"]
    if col is not None:
        fill = etree.SubElement(rpr, _A + "solidFill")
        clr = etree.SubElement(fill, _A + "srgbClr")
        clr.set("val", str(col))
    t = etree.SubElement(new_r, _A + "t")
    t.text = text
    return new_r


def shrink_refs_in_frame(text_frame, *, ref_frac=0.52, ref_color=LIGHT):
    for para in text_frame.paragraphs:
        for run in list(para.runs):
            txt = run.text
            if not txt or "[" not in txt:
                continue
            matches = list(_REF_RE.finditer(txt))
            if not matches:
                continue
            props = _run_props(run)
            run.text = txt[:matches[0].start()]
            anchor = run._r
            for i, m in enumerate(matches):
                anchor = _clone_run_after(anchor, props, m.group(), ref=True,
                                          ref_frac=ref_frac, ref_color=ref_color)
                nxt = matches[i + 1].start() if i + 1 < len(matches) else len(txt)
                between = txt[m.end():nxt]
                if between:
                    anchor = _clone_run_after(anchor, props, between, ref=False)
    return text_frame


# ============================================================
# inject_ref_markers — append [N] at ANCHORS anchors on a slide, then shrink.
# Returns list of (anchor_substr, matched?) for verification.
# ============================================================
def _iter_frames(slide):
    for shp in slide.shapes:
        if shp.has_text_frame:
            yield shp.text_frame


def inject_ref_markers(slide, sid):
    anchors = ANCHORS.get(sid, [])
    report = []
    for ref_nums, anchor in anchors:
        marker = f"[{ref_nums}]"
        placed = False
        for tf in _iter_frames(slide):
            if placed:
                break
            for para in tf.paragraphs:
                if placed:
                    break
                for run in para.runs:
                    if anchor in run.text and marker not in run.text:
                        run.text = run.text.replace(anchor, anchor + marker, 1)
                        placed = True
                        break
        report.append((anchor, placed))
    # shrink every frame that now carries a [N]
    for tf in _iter_frames(slide):
        shrink_refs_in_frame(tf)
    return report


# ============================================================
# ref_list / refs_of_slide — bottom clickable numbered source list.
# ============================================================
def ref_list(slide, entries, *, y=7.06, x=0.55, w=12.25, h=0.36, size=8.5,
             color=LIGHT, line_spacing=1.02, tail=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    for i, (num, name, url) in enumerate(entries):
        rm = p.add_run(); rm.text = f"[{num}] "
        rm.font.name = FONT_BODY; rm.font.size = Pt(size)
        rm.font.bold = True; rm.font.italic = True; rm.font.color.rgb = MID
        rn = p.add_run(); rn.text = name
        rn.font.name = FONT_BODY; rn.font.size = Pt(size)
        rn.font.italic = True; rn.font.color.rgb = color
        if url:
            try:
                rn.hyperlink.address = url
            except Exception:
                pass
        if i < len(entries) - 1:
            rs = p.add_run(); rs.text = "   ·   "
            rs.font.name = FONT_BODY; rs.font.size = Pt(size)
            rs.font.italic = True; rs.font.color.rgb = color
    if tail:
        rt = p.add_run(); rt.text = "   ·   " + tail
        rt.font.name = FONT_BODY; rt.font.size = Pt(size)
        rt.font.italic = True; rt.font.color.rgb = SLATE
    return tb


def _resolve_refs(sid):
    out = []
    for entry in SLIDE_REFS.get(sid, []):
        num, name, urlkey, gloss = entry[0], entry[1], entry[2], entry[3]
        volatile = len(entry) > 4 and entry[4]
        url = URLS.get(urlkey, "") if urlkey else ""
        out.append((num, name, url, gloss, volatile))
    return out


def refs_of_slide(slide, sid, *, y=7.06, x=0.55, w=12.25, tail=None):
    resolved = _resolve_refs(sid)
    if not resolved:
        return None
    entries = [(num, name, url) for (num, name, url, gloss, vol) in resolved]
    n = len(entries)
    size = 8.5 if n <= 3 else (8.0 if n <= 4 else 7.4)
    return ref_list(slide, entries, y=y, x=x, w=w, size=size, tail=tail)


# ============================================================
# notes «Источники:» block
# ============================================================
def notes_sources_block(sid):
    resolved = _resolve_refs(sid)
    if not resolved:
        return ""
    lines = ["Источники:"]
    for (num, name, url, gloss, vol) in resolved:
        if url:
            vfy = " [VFY-day-of]" if vol else ""
            lines.append(f"[{num}] {name} — {gloss}. {url}{vfy}")
        else:
            # course-internal / unconfirmed → [VFY], no canonical URL on record
            lines.append(f"[{num}] {name} — {gloss}. [VFY: не подтверждён canonical-URL, "
                         f"подавать как данные независимого реестра live-eval, не как первоисточник]")
    return "\n".join(lines)


# ============================================================
# page number — muted «N / TOTAL» bottom-right
# ============================================================
def page_number(slide, n, total=None, *, color=SLATE):
    txt = f"{n} / {total}" if total else str(n)
    tb = slide.shapes.add_textbox(Inches(12.33), Inches(7.16), Inches(0.95),
                                  Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; p.line_spacing = 1.0
    r = p.add_run(); r.text = txt
    r.font.name = FONT_BODY; r.font.size = Pt(10); r.font.italic = True
    r.font.color.rgb = color
    return tb
