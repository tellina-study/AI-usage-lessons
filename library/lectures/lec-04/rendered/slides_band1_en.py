"""Лекция 4 v4 — Band 1 (Раздел 0 s01–s07 + Раздел 1 s08–s10).

Methodology-first re-spine. Each sNN(p) builds one slide from slides/sNN*.md
(visible content + visual_brief). Palette Ocean LOCKED, motif «Ocean rounded
box», Gold ≥1×/slide. Visible clean source lines (author/method/year/domain)
on every recommendation/failure slide.
"""
from _helpers_en import (
    blank, set_slide_bg, text_box, text_runs, ocean_box, filled_rect,
    right_arrow, circle, chip, connector, add_image, icon, slide_title,
    gold_callout, teal_callout, footer, src, speaker_notes, load_notes, notes_with_sources, refs_of_slide,
    roadmap_bar, build_section_divider, NAV, ref_list, refs_of, link_run, URLS,
    DEEP, MID, LIGHT, TEAL, SURFACE, WHITE, GOLD, SLATE, COVER_OUTLINE,
    GOLD_TINT, TEAL_TINT, SOFT_GREY, MID_TINT, GOLD_TINT as GT,
    ICONS, CHARTS, ASSETS,
)
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

SCR = ASSETS / "screenshots"


# ============================================================
# s01 — hero_cover (METR perception gap)
# ============================================================
def s01(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(
        s, "16 experts were sure AI sped them up — and got the sign wrong",
        size=23, w=10.5, h=0.98, y=0.36)
    icon(s, "gauge", 12.15, 0.40, 0.78, "mid")

    # HERO — real METR chart, left ~46% area
    hx, hy, hw, hh = 0.55, 1.55, 6.25, 4.15
    ocean_box(s, hx - 0.05, hy - 0.05, hw + 0.10, hh + 0.10,
              fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, SCR / "s01-hero.jpg", hx + 0.10, hy + 0.10,
              hw - 0.20, hh - 0.55)
    text_box(s, x=hx + 0.15, y=hy + hh - 0.42, w=hw - 0.30, h=0.34,
             text="Predicted −24% · believed −20% · actual +19% time",
             size=12.5, italic=True, color=MID, align=PP_ALIGN.CENTER)

    # Setup box top-right
    box_x = 7.10
    ocean_box(s, box_x, 1.55, 5.70, 1.55)
    text_runs(s, box_x + 0.28, 1.70, 5.15, 1.30, [
        {"text": "METR experiment, first half of 2025",
         "size": 15, "bold": True, "color": DEEP},
        {"text": "16 experienced open-source developers · 246 real tasks "
                 "in their own familiar repositories; measured real time, not perception",
         "size": 12.5, "color": SLATE, "newpara": True, "space_before": 6},
    ])

    # Three numbers box
    ocean_box(s, box_x, 3.24, 5.70, 2.42)
    text_runs(s, box_x + 0.28, 3.38, 5.15, 2.18, [
        {"text": "Three numbers about one and the same thing", "size": 14.5, "bold": True,
         "color": MID},
        {"text": "Predicted before the experiment: AI will speed up by −24%",
         "size": 12.5, "color": DEEP, "newpara": True, "space_before": 5},
        {"text": "Believed after working: sped up by roughly −20%",
         "size": 12.5, "color": DEEP, "newpara": True, "space_before": 2},
        {"text": "Measured actual: with AI, tasks took +19% longer",
         "size": 12.5, "bold": True, "color": DEEP, "newpara": True,
         "space_before": 2},
        {"text": "— METR RCT, n=16 experienced OSS developers, 246 tasks, 2025",
         "size": 9, "italic": True, "color": LIGHT, "newpara": True,
         "space_before": 4},
    ])

    # Gold callout bottom (spanning under hero)
    gold_callout(
        s, 0.55, 5.86, 12.25, 0.66,
        "The professionals got the magnitude right but the sign wrong. \"It "
        "feels like the tool helps\" is a hypothesis, not data; reliability comes "
        "from discipline with built-in verification, not from perception. [1]",
        size=13, bold=True)
    refs_of_slide(s, "s01")
    notes_with_sources(s, "s01")
    return s


# ============================================================
# s02 — cover + roadmap
# ============================================================
def s02(p):
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # giant decorative "04" outline, left
    text_box(s, x=-0.10, y=0.55, w=6.2, h=4.6, text="04",
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # decorative chain hint of git-loop artefacts around the number
    chain = ["spec", "ADR", "plan", "PR", "incident"]
    cy = 5.05
    cx = 0.65
    for i, w in enumerate(chain):
        chip(s, cx, cy, 1.02, 0.40, w, fill=(GOLD if i == 0 else MID),
             color=(DEEP if i == 0 else WHITE), size=11.5)
        cx += 1.02
        if i < len(chain) - 1:
            text_box(s, x=cx - 0.02, y=cy + 0.02, w=0.22, h=0.36, text="→",
                     size=15, bold=True, color=LIGHT, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += 0.20

    # title block, right
    text_box(s, x=6.55, y=1.55, w=6.35, h=0.5, text="LECTURE 4",
             size=20, bold=True, color=TEAL)
    text_box(s, x=6.55, y=2.10, w=6.40, h=2.0,
             text="AI across the software development lifecycle",
             size=36, bold=True, color=DEEP, line_spacing=1.05)
    text_box(s, x=6.58, y=4.15, w=6.30, h=0.6,
             text="Course \"Deliberate use of AI\" · software developers",
             size=15, italic=True, color=LIGHT)
    gold_callout(
        s, 6.55, 4.85, 6.35, 0.78,
        "Reliability comes not from the tool but from engineering discipline by "
        "phase: which artifact and which check each one needs.",
        size=13, bold=True)
    roadmap_bar(s, 0, y=6.55)
    notes_with_sources(s, "s02")
    return s


# ============================================================
# s03 — bridge from Module 1 (4 carry-over cards)
# ============================================================
def s03(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(
        s, "We stand on the foundation of Module 1 — take one industry and break it down by phase",
        size=23, w=12.0, h=0.9)

    cards = [
        ("layers", "From Lecture 1",
         "The layered picture \"model → chat → agent → application\"; a prompt as "
         "\"role + task + context\". We use it as a vocabulary."),
        ("braces", "From Lecture 2",
         "Why AI produces \"almost right\" text: the answer is generated "
         "token by token, so plausible ≠ correct. In one phrase."),
        ("route", "From Lecture 3",
         "The complexity ladder (\"stay on the lowest rung\") + the criterion "
         "\"when not AI at all\" — a deterministic task is solved by ordinary code."),
        ("refresh-cw", "From Lecture 3",
         "The agent loop plan → act → check → iterate (4 points of failure) + prompt "
         "injection (defense is architectural, not by filtering)."),
    ]
    cw, ch, gap = 5.95, 1.62, 0.35
    xs = [0.55, 0.55 + cw + gap]
    ys = [1.55, 1.55 + ch + 0.28]
    for i, (ic, head, body) in enumerate(cards):
        x = xs[i % 2]
        y = ys[i // 2]
        ocean_box(s, x, y, cw, ch)
        icon(s, ic, x + 0.24, y + 0.24, 0.62, "mid")
        text_box(s, x=x + 1.06, y=y + 0.20, w=cw - 1.30, h=0.36, text=head,
                 size=14.5, bold=True, color=MID)
        text_box(s, x=x + 1.06, y=y + 0.58, w=cw - 1.30, h=0.95, text=body,
                 size=12, color=DEEP, line_spacing=1.12)

    gold_callout(
        s, 0.55, 5.62, 12.25, 0.98,
        "We carry this over intact and don't re-explain it — we only apply it to software development. "
        "Module 1 gave the apparatus for choosing an architecture; Lecture 4 shows which "
        "discipline makes AI in one industry reliable by phase — and where it breaks without it.",
        size=13, bold=True)
    notes_with_sources(s, "s03")
    return s


# ============================================================
# s04 — central question (contrast two framings)
# ============================================================
def s04(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "The lecture's question is about discipline, not about the tool",
                size=25, w=12.0, h=0.85)

    # central question box
    ocean_box(s, 0.55, 1.50, 12.25, 1.55)
    text_runs(s, 0.85, 1.66, 11.65, 1.30, [
        {"text": "AI writes code ever better — but what makes AI development reliable?",
         "size": 21, "bold": True, "color": DEEP},
        {"text": "Not the tool, but engineering discipline by phase: which "
                 "human-owned artifact and which check each phase needs, where "
                 "the methods of different players converge — and what is not "
                 "delegated to the tool in any phase. [1]",
         "size": 13.5, "color": SLATE, "newpara": True, "space_before": 8,
         "line_spacing": 1.15},
    ])

    # contrast two frames
    fy = 3.28
    # struck-through weak frame (light, teal-tint)
    filled_rect(s, 0.55, fy, 6.0, 1.55, SOFT_GREY, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.06)
    icon(s, "circle-slash", 0.78, fy + 0.22, 0.5, "light")
    text_box(s, x=1.42, y=fy + 0.22, w=4.9, h=0.5,
             text="Which tool is best?", size=15, bold=True, color=SLATE)
    text_box(s, x=0.80, y=fy + 0.80, w=5.5, h=0.65,
             text="Goes stale in a quarter; \"best\" can't be named without "
                  "\"for what, in which phase, in which mode\".",
             size=12, color=SLATE, line_spacing=1.12)
    # strong frame (gold)
    filled_rect(s, 6.80, fy, 6.0, 1.55, GOLD_TINT, stroke=GOLD, stroke_pt=1.8,
                radius=True, radius_adj=0.06)
    icon(s, "check-check", 7.03, fy + 0.22, 0.5, "gold")
    text_box(s, x=7.67, y=fy + 0.20, w=4.9, h=0.55,
             text="Which practice is warranted in the phase?", size=15, bold=True,
             color=DEEP)
    text_box(s, x=7.05, y=fy + 0.80, w=5.5, h=0.65,
             text="Which artifact does it end in? Where is a human mandatory? — "
                  "stable for years.",
             size=12, color=DEEP, line_spacing=1.12)

    gold_callout(
        s, 0.55, 5.10, 12.25, 0.80,
        "The tool is secondary, the practice primary. When answering, an engineer names not "
        "a logo but the phase's appropriate practice, its artifact, and the point of mandatory "
        "human control — knowledge that will outlive the swap of any vendor. [2]",
        size=13, bold=True)
    refs_of_slide(s, "s04")
    notes_with_sources(s, "s04")
    return s


# ============================================================
# s05 — NEW foundations slide (ПРАВКА 3, #266b) — эта лекция = сводка
# практик: A «современные (лидеры)» + B «классика». Каждый пункт —
# кликабельная ссылка на канонический источник. Вставлен ДО keystone.
# ============================================================
def s05f(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(
        s, "This lecture is a digest of practices: modern approaches of the leaders + time-tested classics",
        size=20, w=12.4, h=0.82)

    # two columns of practice lists
    colw = 6.05
    gap = 0.15
    lx = 0.55
    rx = lx + colw + gap
    top = 1.46
    boxh = 4.62

    # --- Column A: modern practices (leaders) ---
    ocean_box(s, lx, top, colw, boxh, fill=SURFACE, stroke=MID, stroke_pt=1.6)
    icon(s, "sliders-horizontal", lx + 0.24, top + 0.18, 0.46, "mid")
    text_box(s, x=lx + 0.84, y=top + 0.20, w=colw - 1.0, h=0.40,
             text="A. Modern practices (leaders)", size=13.5, bold=True,
             color=MID)
    modern = [
        ("1", "Anthropic — Claude Code / AI-Native SDLC", "anthropic_playbook",
         "agentic git loop: each stage commits an artifact; human at the gates"),
        ("2", "OpenAI — Model Spec", "model_spec",
         "spec-as-contract; clause = example prompt = test"),
        ("3", "GitHub — Spec Kit", "spec_kit",
         "\"intent is the source of truth\"; small verifiable tasks"),
        ("4", "Google — DORA 2025", "dora_2025",
         "\"AI amplifies what is already there\"; seven capabilities"),
        ("5", "Thoughtworks — Exploring Gen AI", "fowler_genai",
         "the assistant suggests, the developer owns; harness engineering"),
        ("6", "Willison — Vibe engineering", "willison_vibe_eng",
         "the disciplines an LLM rewards: tests, plans, reviews"),
    ]
    _foundation_list(s, modern, lx + 0.26, top + 0.72, colw - 0.52)

    # --- Column B: time-tested classics ---
    ocean_box(s, rx, top, colw, boxh, fill=SURFACE, stroke=LIGHT, stroke_pt=1.6)
    icon(s, "graduation-cap", rx + 0.24, top + 0.18, 0.46, "teal")
    text_box(s, x=rx + 0.84, y=top + 0.20, w=colw - 1.0, h=0.40,
             text="B. Time-tested classics", size=13.5, bold=True,
             color=TEAL)
    classics = [
        ("7", "Brooks — No Silver Bullet", "brooks",
         "essential vs accidental complexity; \"what to build\" is the human's"),
        ("8", "Beck — TDD: By Example", "beck_tdd",
         "red-green-refactor; test-as-specification"),
        ("9", "Nygard — ADR", "nygard_adr",
         "immutable \"why\" records under version control"),
        ("10", "Ford/Parsons — Evolutionary Architectures", "evol_arch",
         "fitness functions: \"fitness\" is objective and automatic"),
        ("11", "Fowler — Refactoring", "fowler_refactoring",
         "the discipline of small verifiable changes"),
        ("12", "Brown — C4", "c4",
         "architecture-as-code: a textual, diffable model for AI"),
    ]
    _foundation_list(s, classics, rx + 0.26, top + 0.72, colw - 0.52,
                     accent=TEAL)

    gold_callout(
        s, 0.55, 6.24, 12.25, 0.52,
        "Tools will change over quarters — these practices are stable: they rest "
        "on the nature of complexity, not on product maturity. All names are "
        "clickable links to the primary sources.",
        size=12, bold=True, align=PP_ALIGN.CENTER)
    notes_with_sources(s, "s05")
    return s


def _foundation_list(s, items, x, y0, w, *, accent=MID):
    """Render a foundation list column: [N] clickable name + 1-line body."""
    row_h = 0.635
    for i, (num, name, urlkey, body) in enumerate(items):
        y = y0 + i * row_h
        # [N] marker + clickable name on line 1
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.30))
        tf = tb.text_frame
        tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
        tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
        tf.word_wrap = True
        pgraph = tf.paragraphs[0]
        pgraph.line_spacing = 1.0
        link_run(pgraph, f"[{num}] ", "", size=10.5, color=GOLD, bold=True)
        link_run(pgraph, name, URLS.get(urlkey, ""), size=10.5, color=DEEP,
                 bold=True)
        # body on line 2 (muted)
        text_box(s, x=x + 0.10, y=y + 0.29, w=w - 0.10, h=0.30, text=body,
                 size=9.5, italic=True, color=SLATE, line_spacing=1.0)


# ============================================================
# s06 — KEYSTONE — цикл ФАЗ лекции (ПРАВКА 2, #265)
# Совпадает с дивайдерами/роадмапом 0–7: Требования → Архитектура →
# Реализация → Тестирование → Ревью+Безопасность → Доставка/Эксплуатация →
# Документация → (обратно к Требованиям). Каждый узел-фаза подписан своим
# ЧЕЛОВЕКО-ВЛАДЕЕМЫМ артефактом.
# ============================================================
def s06k(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(
        s, "Discipline is a cycle of phases, and each phase is owned by a human through its artifact",
        size=21, w=12.3, h=0.85)

    # ---- phase cycle: 7 nodes in two rows, forward arrows, gold return ----
    # Node = (icon, phase-name, human-owned artefact, strength)
    #   strong=MID stroke; thin=LIGHT stroke; bright=GOLD (docs)
    nodes = [
        ("file-code", "1 Requirements", "requirements.md", "strong"),
        ("gavel", "2 Architecture", "ADR", "thin"),
        ("code", "3 Implementation", "PR / code", "strong"),
        ("flask-conical", "4 Testing", "test-as-spec", "strong"),
        ("shield-check", "5 Review+Sec.", "review gate", "strong"),
        ("git-merge", "6 Delivery/Ops", "release gate", "thin"),
        ("lightbulb", "7 Documentation", "docs-as-context", "bright"),
    ]
    # layout: 4 on top row, 3 on bottom row
    row1 = nodes[:4]
    row2 = nodes[4:]
    x0 = 0.55
    total = 12.25
    gap = 0.30
    nw = (total - gap * 3) / 4        # 4 columns ~2.79
    nh = 1.30
    y1 = 1.50
    y2 = 3.06

    def draw_node(ix, iy, ic, name, art, strength):
        if strength == "bright":
            fill, stroke, spt, av = GOLD_TINT, GOLD, 1.8, "gold"
        elif strength == "thin":
            fill, stroke, spt, av = SURFACE, LIGHT, 1.2, "light"
        else:
            fill, stroke, spt, av = SURFACE, MID, 1.6, "mid"
        ocean_box(s, ix, iy, nw, nh, fill=fill, stroke=stroke, stroke_pt=spt)
        icon(s, ic, ix + 0.20, iy + 0.18, 0.50, av)
        text_box(s, x=ix + 0.80, y=iy + 0.20, w=nw - 0.94, h=0.46, text=name,
                 size=13, bold=True, color=DEEP, line_spacing=0.98)
        # artefact chip (human-owned)
        chip(s, ix + 0.20, iy + 0.86, nw - 0.40, 0.40, art,
             fill=(GOLD if strength == "bright" else MID),
             color=(DEEP if strength == "bright" else WHITE), size=10.5)

    centers1 = []
    for i, (ic, name, art, strn) in enumerate(row1):
        x = x0 + i * (nw + gap)
        draw_node(x, y1, ic, name, art, strn)
        centers1.append((x, x + nw))
        if i < 3:
            right_arrow(s, x + nw + 0.02, y1 + nh / 2 - 0.13, gap - 0.06, 0.26,
                        fill=LIGHT)
    # connector row1 → row2 (down on the right, curve implied by arrow)
    last1_cx = centers1[-1][0] + nw / 2
    text_box(s, x=last1_cx - 0.30, y=y1 + nh + 0.01, w=0.6, h=0.28, text="▼",
             size=13, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)

    centers2 = []
    # row2 laid out left→right but reading continues 5,6,7
    for i, (ic, name, art, strn) in enumerate(row2):
        x = x0 + i * (nw + gap)
        draw_node(x, y2, ic, name, art, strn)
        centers2.append((x, x + nw))
        if i < len(row2) - 1:
            right_arrow(s, x + nw + 0.02, y2 + nh / 2 - 0.13, gap - 0.06, 0.26,
                        fill=LIGHT)
    # gold RETURN arrow: Документация(7) → Требования(1), closes the cycle.
    # Routed compactly just below row2 so it doesn't collide with the info row.
    ret_x = x0 + (len(row2) - 1) * (nw + gap) + nw / 2   # centre of node 7
    first_cx = x0 + nw / 2                               # centre of node 1
    ry = y2 + nh + 0.24
    connector(s, ret_x, y2 + nh, ret_x, ry, color=GOLD, width=2.6)
    connector(s, ret_x, ry, first_cx, ry, color=GOLD, width=2.6)
    connector(s, first_cx, ry, first_cx, y1 + nh, color=GOLD, width=2.6)
    text_box(s, x=first_cx - 0.30, y=y1 + nh - 0.02, w=0.6, h=0.28, text="▲",
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # label placed in the open space to the RIGHT of node 7, above the return leg
    text_box(s, x=ret_x + 0.20, y=ry - 0.30, w=3.9, h=0.56,
             text="the cycle closes: an ops lesson → a new requirement",
             size=10.5, italic=True, bold=True, color=GOLD,
             align=PP_ALIGN.LEFT, line_spacing=1.05)

    # bottom-left: phase definition + AI-inside/human-owns
    dy = ry + 0.36
    ocean_box(s, 0.55, dy, 7.35, 1.12)
    text_runs(s, 0.80, dy + 0.10, 6.85, 0.96, [
        {"text": "A phase is a stage with its own input / output / artifact.",
         "size": 12.5, "bold": True, "color": MID},
        {"text": "AI works INSIDE the nodes (drafts requirements, proposes "
                 "an ADR, writes code in a PR), but each node is owned by a human — who reads, "
                 "edits, accepts, is accountable [1].",
         "size": 11, "color": DEEP, "newpara": True, "space_before": 4,
         "line_spacing": 1.10},
    ])
    # bottom-right: uneven maturity
    filled_rect(s, 8.10, dy, 4.70, 1.12, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.06)
    text_runs(s, 8.34, dy + 0.08, 4.25, 0.98, [
        {"text": "The practices mature unevenly", "size": 11.5, "bold": True,
         "color": DEEP},
        {"text": "strong: requirements, implementation, testing, review · "
                 "thin: architecture, delivery, operations · bright spot: "
                 "documentation.",
         "size": 10.5, "color": SLATE, "newpara": True, "space_before": 3,
         "line_spacing": 1.08},
    ])

    gold_callout(
        s, 0.55, dy + 1.22, 12.25, 0.44,
        "Anthropic [1], OpenAI [2], DORA [3], and Thoughtworks [4] independently "
        "converged on this phase skeleton — so what we have is a method, not a fashion.",
        size=11.5, bold=True, align=PP_ALIGN.CENTER)

    refs_of_slide(s, "s06")
    notes_with_sources(s, "s06")
    return s


# ============================================================
# s06 — autonomy ladder as demoted lens
# ============================================================
def s06(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Autonomy is a property of the mode, not the brand",
                size=25, w=12.0, h=0.85)

    # left: ladder A->D
    lx, lw = 0.55, 5.85
    ocean_box(s, lx, 1.55, lw, 4.05)
    text_box(s, x=lx + 0.26, y=1.70, w=lw - 0.52, h=0.44,
             text="The autonomy ladder — a supporting lens",
             size=13.5, bold=True, color=MID)
    steps = [
        ("A", "autocomplete", "finishes the line; the human accepts each one"),
        ("B", "small tasks", "a function/fix in a dialog; the human reviews after"),
        ("C", "coding agent", "plans, edits files, runs tests; the human does the merge"),
        ("D", "orchestrator", "takes a task from the tracker → PR; the human is the prod gate"),
    ]
    sy = 2.24
    for i, (lab, name, desc) in enumerate(steps):
        y = sy + i * 0.80
        circle(s, lx + 0.28, y, 0.52, MID)
        text_box(s, x=lx + 0.28, y=y, w=0.52, h=0.52, text=lab, size=18,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=lx + 0.98, y=y - 0.02, w=lw - 1.30, h=0.34, text=name,
                 size=13.5, bold=True, color=DEEP)
        text_box(s, x=lx + 0.98, y=y + 0.32, w=lw - 1.24, h=0.42, text=desc,
                 size=11, color=SLATE, line_spacing=1.02)
    icon(s, "arrow-right-left", lx + lw - 0.78, 1.68, 0.5, "teal")

    # right: mode != brand
    rx, rw = 6.70, 6.10
    ocean_box(s, rx, 1.55, rw, 1.92)
    text_box(s, x=rx + 0.26, y=1.68, w=rw - 0.52, h=0.40,
             text="Mode ≠ brand — one product lives on several rungs at once [2]",
             size=13.5, bold=True, color=MID, line_spacing=1.05)
    text_runs(s, rx + 0.26, 2.42, rw - 0.52, 0.95, [
        {"text": "Copilot — A, B, C and D · Cursor — Tab (A), Cmd-K (B), "
                 "Composer (C) · Claude Code — C, rising to D",
         "size": 12.5, "color": DEEP, "line_spacing": 1.18},
    ])
    # two boundaries
    filled_rect(s, rx, 3.62, rw, 0.90, TEAL_TINT, stroke=TEAL, stroke_pt=1.4,
                radius=True, radius_adj=0.07)
    text_runs(s, rx + 0.26, 3.72, rw - 0.52, 0.72, [
        {"text": "B ↔ C  ", "size": 13, "bold": True, "color": TEAL},
        {"text": "does it iterate and run tests without you? (yes → C)",
         "size": 12, "color": DEEP},
        {"text": "C ↔ D  ", "size": 13, "bold": True, "color": TEAL,
         "newpara": True, "space_before": 3},
        {"text": "where the task comes from, where it lands (tracker → PR → D)",
         "size": 12, "color": DEEP},
    ])

    gold_callout(
        s, 6.70, 4.66, 6.10, 0.94,
        "The higher the autonomy, the stricter the criterion \"a human is mandatory here\" and "
        "the tighter the harness. \"We use Copilot\" reports neither the level nor the phase "
        "— name the mode and the phase, not the logo. [1]",
        size=12.5, bold=True)
    refs_of_slide(s, "s07")
    notes_with_sources(s, "s07")
    return s


# ============================================================
# s07 — tool-secondary thesis (in-bucket)
# ============================================================
def s07(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "The method decides, the tool executes",
                size=26, w=12.0, h=0.85)

    ocean_box(s, 0.55, 1.50, 12.25, 1.05)
    text_box(s, x=0.85, y=1.62, w=11.65, h=0.85,
             text="A methodological practice is the decision about which artifact "
                  "to produce, in what order, with what check, and who is accountable. "
                  "The tool is the executor: interchangeable and secondary.",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)

    # three pillars
    pillars = [
        ("scale", "Volatility",
         "Tools and their \"leadership\" change over quarters; practices are stable "
         "for years — they rest on the nature of complexity (Brooks [1]), not on product maturity."),
        ("sliders-horizontal", "The DORA multiplier",
         "\"AI amplifies what is already there\" [2]: build the discipline first, not "
         "pick the tool — a tool without discipline multiplies chaos."),
        ("shield-check", "Accountability",
         "The tool is not accountable for consequences, the human is; accountability "
         "materializes in artifacts and gates, that is, in the practice."),
    ]
    cw, gap = 3.97, 0.17
    x0 = 0.55
    py = 2.72
    for i, (ic, head, body) in enumerate(pillars):
        x = x0 + i * (cw + gap)
        ocean_box(s, x, py, cw, 1.72)
        icon(s, ic, x + 0.24, py + 0.22, 0.54, "teal")
        text_box(s, x=x + 0.90, y=py + 0.26, w=cw - 1.10, h=0.42, text=head,
                 size=14, bold=True, color=MID)
        text_box(s, x=x + 0.24, y=py + 0.80, w=cw - 0.48, h=0.86, text=body,
                 size=11, color=DEEP, line_spacing=1.12)

    # failure strip
    filled_rect(s, 0.55, 4.62, 12.25, 0.98, SOFT_GREY, stroke=LIGHT,
                stroke_pt=1.2, radius=True, radius_adj=0.05)
    icon(s, "triangle-alert", 0.80, 4.82, 0.52, "light")
    text_box(s, x=1.48, y=4.72, w=11.05, h=0.82,
             text="\"We adopted an AI tool\" ≠ \"we adopted an AI discipline\". "
                  "The tool is there, the practice is not — and the DORA multiplier works "
                  "the wrong way. Behind this stand all the failures of the lecture: prompt-and-pray, "
                  "poisoned context, the 70% problem, the Replit incident — everywhere "
                  "the tool applied without the practice.",
             size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)

    gold_callout(
        s, 0.55, 5.72, 12.25, 0.62,
        "Method-first order: first the phase's practice (artifact, gate, who is "
        "accountable), then the tool to fit it. The tool is chosen to fit the "
        "discipline, not the discipline to fit the tool. [3]",
        size=13, bold=True, align=PP_ALIGN.CENTER)
    refs_of_slide(s, "s08")
    notes_with_sources(s, "s08")
    return s


# ============================================================
# s08 — section divider Раздел 1 (Требования)
# ============================================================
def s08(p):
    return build_section_divider(
        p, here_idx=1,
        subtitle="Requirements — the first artifact, before any code",
        bridge="The cycle's first phase turns a vague intent into requirements. "
               "What leads here is not the tool but the discipline of requirements "
               "(spec-driven): between intent and code a "
               "human-reviewable artifact is placed, and the decision \"what exactly we need\" "
               "stays human.",
        sid="s09")


# ============================================================
# s09 — spec-driven practice (git tree of spec files + 3 voices)
# ============================================================
def s09(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "The first development artifact is requirements, not code",
                size=25, w=12.0, h=0.85)

    # left: git tree of requirement files (main visual) — RU names (ПРАВКА 5)
    lx, lw = 0.55, 5.55
    ocean_box(s, lx, 1.55, lw, 4.05)
    icon(s, "git-branch", lx + 0.24, 1.70, 0.52, "mid")
    text_box(s, x=lx + 0.90, y=1.76, w=lw - 1.10, h=0.40,
             text="Repository: requirements are versioned next to the code [3]",
             size=13, bold=True, color=MID, line_spacing=1.0)
    tree = [
        ("file-code", "requirements.md (spec.md)", "what the system must do", GOLD),
        ("file-code", "design.md", "constraints and decisions", GOLD),
        ("file-code", "tasks.md", "decomposition into tasks", GOLD),
        ("code", "src/…", "code — generated from requirements", MID),
    ]
    ty = 2.36
    for i, (ic, name, desc, col) in enumerate(tree):
        y = ty + i * 0.72
        indent = lx + 0.40
        icon(s, ic, indent, y, 0.42, "gold" if col == GOLD else "mid")
        text_box(s, x=indent + 0.58, y=y - 0.02, w=lw - 1.10, h=0.34,
                 text=name, size=13, bold=True, color=DEEP,
                 font="DejaVu Sans Mono")
        text_box(s, x=indent + 0.58, y=y + 0.32, w=lw - 1.10, h=0.32,
                 text=desc, size=11, italic=True, color=SLATE)
    text_box(s, x=lx + 0.26, y=5.16, w=lw - 0.52, h=0.40,
             text="The order is enforced: requirements → design → tasks "
                  "(Kiro / Spec-Kit [3]). The costliest errors are in this phase.",
             size=10, italic=True, bold=True, color=MID, line_spacing=1.02)

    # right: three voices
    rx, rw = 6.30, 6.50
    voices = [
        ("OpenAI Model Spec [2]",
         "living, versioned Markdown: requirements are a durable, reviewable, "
         "diffable artifact, not fleeting prompts."),
        ("Sean Grove, OpenAI · \"The New Code\" [5]",
         "\"the source requirements — that is the valuable artifact\"; code is \"structured "
         "communication\"."),
        ("Martin Fowler [4]",
         "the bottleneck of AI development is intent: the model writes code ever better, but "
         "expressing what to build is still hard."),
    ]
    vy = 1.55
    for i, (head, body) in enumerate(voices):
        y = vy + i * 0.98
        ocean_box(s, rx, y, rw, 0.90)
        text_box(s, x=rx + 0.24, y=y + 0.09, w=rw - 0.48, h=0.30, text=head,
                 size=12.5, bold=True, color=MID)
        text_box(s, x=rx + 0.24, y=y + 0.38, w=rw - 0.48, h=0.50, text=body,
                 size=10.5, color=DEEP, line_spacing=1.06)
    # caveat about ~10-20%
    filled_rect(s, rx, 4.53, rw, 0.60, SOFT_GREY, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.08)
    text_box(s, x=rx + 0.24, y=4.60, w=rw - 0.48, h=0.46,
             text="Grove's \"code ~10–20% of the value\" is a rhetorical provocation, not "
                  "an empirical measurement: its job is to shift the focus to requirements.",
             size=10.5, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)

    gold_callout(
        s, 0.55, 5.72, 12.25, 0.58,
        "AI is strong at structuring and completing intent; the intent itself — "
        "what the system needs — is decided by the human [1].",
        size=13, bold=True, align=PP_ALIGN.CENTER)

    refs_of_slide(s, "s10")
    notes_with_sources(s, "s10")
    return s


# ============================================================
# s10 (display s11) — как вести требования: СТРУКТУРА + ПРОЦЕСС (ПРАВКА 6)
# ============================================================
def s10(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(
        s, "How to run requirements: structure (how to write) + process (how to maintain)",
        size=21, w=12.4, h=0.82)

    colw = 6.05
    gap = 0.15
    lx = 0.55
    rx = lx + colw + gap
    top = 1.44
    boxh = 4.02

    # --- LEFT: STRUCTURE ---
    ocean_box(s, lx, top, colw, boxh, fill=SURFACE, stroke=MID, stroke_pt=1.6)
    icon(s, "list-checks", lx + 0.24, top + 0.16, 0.44, "mid")
    text_box(s, x=lx + 0.82, y=top + 0.18, w=colw - 1.0, h=0.36,
             text="STRUCTURE — how to write requirements", size=13, bold=True,
             color=MID)
    struct = [
        ("Stories + acceptance criteria [3]",
         "\"As a <role>, I want <goal>, so that <benefit>\" + verifiable criteria for "
         "each story."),
        ("EARS notation [7]",
         "\"WHEN <trigger>, the system SHALL <response>\" (Mavin 2009) — removes "
         "\"should/may\", makes a requirement testable."),
        ("Functional vs non-functional [6]",
         "behavior separate from characteristics (latency / cost / security); NFRs "
         "are enforced by fitness functions."),
        ("requirements → design → tasks [3]",
         "the enforced order of 3 files (Kiro / Spec-Kit); DoD — small "
         "independently testable units."),
    ]
    sy = top + 0.62
    for i, (head, body) in enumerate(struct):
        y = sy + i * 0.82
        text_box(s, x=lx + 0.26, y=y, w=colw - 0.52, h=0.30, text=f"• {head}",
                 size=11, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=lx + 0.42, y=y + 0.28, w=colw - 0.68, h=0.50, text=body,
                 size=10, color=SLATE, line_spacing=1.06)

    # --- RIGHT: PROCESS ---
    ocean_box(s, rx, top, colw, boxh, fill=SURFACE, stroke=LIGHT, stroke_pt=1.6)
    icon(s, "refresh-cw", rx + 0.24, top + 0.16, 0.44, "teal")
    text_box(s, x=rx + 0.82, y=top + 0.18, w=colw - 1.0, h=0.36,
             text="PROCESS — how to maintain requirements", size=13, bold=True,
             color=TEAL)
    proc = [
        ("Elicitation: the interrogatory LLM [4]",
         "the model ASKS questions (Fowler's \"Interrogatory LLM\"), surfacing "
         "unstated assumptions — instead of \"prompt-and-pray\"."),
        ("Review and sign-off BEFORE code [1]",
         "requirements are reviewed and signed off by a human before generation; "
         "accept/reject = \"the merge\"."),
        ("Versioning next to the code [3]",
         "requirements are diffable Markdown in the repository, not in a wiki / chat; "
         "a durable artifact, not a fleeting prompt."),
        ("Syncing with change [9]",
         "keep current like an ADR; the human owns \"what to build\", AI owns "
         "structure and completeness."),
    ]
    for i, (head, body) in enumerate(proc):
        y = sy + i * 0.82
        text_box(s, x=rx + 0.26, y=y, w=colw - 0.52, h=0.30, text=f"• {head}",
                 size=11, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=rx + 0.42, y=y + 0.28, w=colw - 0.68, h=0.50, text=body,
                 size=10, color=SLATE, line_spacing=1.06)

    # judgment plate
    gold_callout(
        s, 0.55, 5.62, 12.25, 0.78,
        "Durable pattern: EARS + decomposition + requirements-as-check + "
        "human sign-off — will outlive any tool. Hype: \"our pipeline of "
        "commands = requirements discipline\". Tools execute, the method leads.",
        size=12.5, bold=True, align=PP_ALIGN.CENTER)

    refs_of_slide(s, "s11")
    notes_with_sources(s, "s11")
    return s
