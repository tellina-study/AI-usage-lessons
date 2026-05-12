"""
Full 29-slide build of Лекции 1 (Phase 5-6 of EPIC #64, issue #69).

Source-of-truth: chapter v2 (status=reviewed, 13268 слов).
Plan: notes/lecture-1-review/final/new-plan-v5-final.md.

Builds upon the v3.6 6-slide pilot (archive-v36-6slide/build_v36.py) — keeps the
same Ocean Gradient + Teal + Gold palette, Ocean rounded box motif,
distinct cover, 16:9 canvas.

Slides:
  s01–s05b — adapted from v3.6 pilot (chapter v2 corrections in metadata,
             content already finalised in v3.6).
  s06–s29 — new, built from chapter v2 content per slide markdown specs.

Canvas: 13.333" × 7.5" (16:9).  Rendered at libreoffice + pdftoppm @150dpi.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)  # very pale gold for callouts
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)  # very pale teal for accents

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ASSETS = Path("/home/levko/AI-usage-lessons/library/lectures/lec-01/rendered/assets")
OUT = Path("/home/levko/AI-usage-lessons/library/lectures/lec-01/rendered/lec-01.pptx")
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    """Create textbox with multiple paragraphs/runs.
    Each run dict: {text, size, bold, italic, color, font, newpara (start a new paragraph before this run)}.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0, radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=16, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def footer_sources(slide, text, *, color=LIGHT, size=12):
    """Footer-style — sources / caveat only, italic, light blue."""
    text_box(slide, x=0.5, y=7.05, w=12.3, h=0.35, text=text,
             size=size, italic=True, color=color, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)


def slide_title(slide, text, *, y=0.45, h=1.0, w=12.3, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.18, align=PP_ALIGN.LEFT):
    """Standard content-slide title (assertion) at unified 26pt bold (down from 28 for fit)."""
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True):
    """Pale-gold-tinted callout box with bold text, single line."""
    box = filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.0,
                      radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.08, w=w - 0.4, h=h - 0.16, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


def speaker_notes(slide, text):
    """Add speaker notes to slide (visible in presentation mode)."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


# ============================================================
# Slide builders — s01-s05b carry over from v3.6 unchanged
# ============================================================

def build_s01(p):
    """s01 — live_demo (carry from v3.6: hook + YOLO mock + Ocean rounded box)."""
    s = blank(p)
    text_box(s, x=0.55, y=0.55, w=5.9, h=2.4,
             text="Идентификация людей в реальном времени — уже с 2023 года на простом ноутбуке.",
             size=28, bold=True, color=DEEP, line_spacing=1.18)
    text_box(s, x=0.55, y=3.25, w=5.9, h=1.4,
             text=("Narrow AI — модель решает одну задачу "
                   "(обнаружение людей в кадре) и больше ничего."),
             size=15, italic=True, color=MID, line_spacing=1.3)
    # Bottom caption
    tb = s.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(5.9), Inches(1.0))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.line_spacing = 1.35
    for cfg in [
        {"text": "На экране — ", "size": 15, "color": DEEP},
        {"text": "YOLOv8", "size": 15, "color": MID, "bold": True},
        {"text": " на CPU ноутбука: ", "size": 15, "color": DEEP},
        {"text": "31 fps", "size": 15, "color": GOLD, "bold": True},
        {"text": ".", "size": 15, "color": DEEP},
    ]:
        r = p1.add_run(); r.text = cfg["text"]
        r.font.name = FONT_BODY; r.font.size = Pt(cfg.get("size", 15))
        r.font.bold = cfg.get("bold", False); r.font.color.rgb = cfg.get("color", DEEP)
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.35
    for cfg in [
        {"text": "Без интернета", "size": 15, "color": TEAL, "bold": True},
        {"text": "  ·  обучена в 2023.", "size": 15, "color": DEEP},
    ]:
        r = p2.add_run(); r.text = cfg["text"]
        r.font.name = FONT_BODY; r.font.size = Pt(cfg.get("size", 15))
        r.font.bold = cfg.get("bold", False); r.font.color.rgb = cfg.get("color", DEEP)

    # Right column — Ocean rounded box framing the YOLO mock screenshot
    box_x, box_y, box_w, box_h = 6.55, 0.55, 6.3, 4.4
    ocean_box(s, box_x, box_y, box_w, box_h)
    pad = 0.18
    img_w = box_w - 2 * pad
    img_h = img_w * 720.0 / 1280.0
    img_x = box_x + pad
    img_y = box_y + (box_h - img_h) / 2.0
    add_image(s, ASSETS / "illustrations/s01-yolo-mock.png", img_x, img_y, img_w, img_h)

    text_box(s, x=box_x, y=box_y + box_h + 0.05, w=box_w, h=0.4,
             text="Кадр модели в момент демо: 2 человека в боксах. YOLOv8 (Ultralytics, 2023).",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, "Live-демо: ноутбук + веб-камера, проектор показывает аудиторию в real-time с bounding-box детекцией. Модель работает локально, без интернета, ~30 fps на CPU. Это narrow AI — рабочая инженерная лошадка. Backup при отказе HDMI/камеры — assets/code/ice-breaker-cv/backup/. LO1 only (LO7 снят v5).")


def build_s02(p):
    """s02 — DISTINCT cover (light tinted bg, decorative «01», 64pt title)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=8.0, y=2.7, w=5.3, h=4.7,
             text="01",
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0)
    text_box(s, x=0.7, y=1.0, w=6.5, h=0.55,
             text="ЛЕКЦИЯ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.72, 1.55, 0.7, 0.05, fill=TEAL)
    text_box(s, x=0.7, y=2.0, w=8.5, h=2.4,
             text="Введение —\nAI вокруг нас",
             size=64, bold=True, color=DEEP, line_spacing=1.05,
             align=PP_ALIGN.LEFT)
    filled_rect(s, 0.7, 5.45, 0.05, 0.55, fill=TEAL)
    text_box(s, x=0.95, y=5.45, w=8.0, h=0.6,
             text="Карта применений AI: где работает, где — нет.",
             size=22, color=MID, italic=False, align=PP_ALIGN.LEFT,
             line_spacing=1.25)
    hero_w = 5.0
    add_image(s, ASSETS / "illustrations/hero-cover-light.png",
              x=8.0, y=0.9, w=hero_w, h=hero_w)
    speaker_notes(s, "Cover-слайд: визуально отличается от content (крупный title 64pt, декоративный «01», tinted background, без motif rounded box). Произносить минимально: «Лекция 1, введение в AI. Сегодня — карта применений: где AI уже работает в индустрии, а где не помогает. Поехали.» Tone: навигационный (показ, не обещание).")


def build_s03(p):
    """s03 — poll questions (carry from v3.6: 2 motif cards, neutral assertion)."""
    s = blank(p)
    text_box(s, x=0.6, y=0.55, w=12.2, h=0.85,
             text="Сначала — ваша оценка, потом — данные.",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT,
             line_spacing=1.15)
    card_y = 2.0
    card_h = 4.4
    card_w = 5.95
    # Q1
    q1_x = 0.55
    ocean_box(s, q1_x, card_y, card_w, card_h)
    add_image(s, ASSETS / "icons/lucide-hand-blue.png",
              x=q1_x + 0.35, y=card_y + 0.35, w=0.95, h=0.95)
    text_box(s, x=q1_x + 1.5, y=card_y + 0.4, w=card_w - 1.6, h=0.35,
             text="Вопрос 1  ·  выберите ОДИН вариант",
             size=14, bold=True, color=MID)
    text_box(s, x=q1_x + 0.4, y=card_y + 1.55, w=card_w - 0.8, h=1.1,
             text="Какой процент россиян использовали AI в 2025?",
             size=22, bold=True, color=DEEP, line_spacing=1.25)
    chip_y = card_y + 3.05
    chip_w = 1.25
    chip_h = 0.55
    gap = 0.13
    options = ["<20%", "20–40%", "40–60%", ">60%"]
    total_w = len(options) * chip_w + (len(options) - 1) * gap
    start_x = q1_x + (card_w - total_w) / 2.0
    for i, opt in enumerate(options):
        chip(s, start_x + i * (chip_w + gap), chip_y, chip_w, chip_h, opt,
             fill=MID, color=WHITE, size=14, bold=True)
    text_box(s, x=q1_x + 0.4, y=chip_y + 0.7, w=card_w - 0.8, h=0.3,
             text="(поднимите руку на одном варианте)",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Q2
    q2_x = q1_x + card_w + 0.4
    ocean_box(s, q2_x, card_y, card_w, card_h)
    add_image(s, ASSETS / "icons/lucide-message-square-blue.png",
              x=q2_x + 0.35, y=card_y + 0.35, w=0.95, h=0.95)
    text_box(s, x=q2_x + 1.5, y=card_y + 0.4, w=card_w - 1.6, h=0.35,
             text="Вопрос 2  ·  можно НЕСКОЛЬКО",
             size=14, bold=True, color=TEAL)
    text_box(s, x=q2_x + 0.4, y=card_y + 1.55, w=card_w - 0.8, h=1.1,
             text="Кто использовал AI сегодня — и для чего?",
             size=22, bold=True, color=DEEP, line_spacing=1.25)
    chip_y2 = card_y + 3.05
    options2 = ["код", "текст", "перевод", "другое"]
    chip_w2 = 1.4
    total_w2 = len(options2) * chip_w2 + (len(options2) - 1) * gap
    start_x2 = q2_x + (card_w - total_w2) / 2.0
    for i, opt in enumerate(options2):
        chip(s, start_x2 + i * (chip_w2 + gap), chip_y2, chip_w2, chip_h, opt,
             fill=WHITE, stroke=TEAL, color=TEAL, size=14, bold=True)
    text_box(s, x=q2_x + 0.4, y=chip_y2 + 0.7, w=card_w - 0.8, h=0.3,
             text="(поднимите руку столько раз, для скольки задач использовали)",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, "Step 1 reveal-последовательности (парный с s04). Цифр и ответов на этом слайде нет сознательно — это шаг 1, раскрытие на s04. Считаем руки по обоим вопросам. Запомнить распределение для s04 (эффект «ваша оценка vs реальность»). Темп: ~45 сек на вопрос. Не комментировать ответы — просто фиксировать руки.")


def build_s04(p):
    """s04 — poll reveal (donut 51% + bar 5 LLM, both multi-select)."""
    s = blank(p)
    text_box(s, x=0.55, y=0.4, w=12.3, h=0.9,
             text="Разница между вашей оценкой и реальностью — карта ваших слепых зон про AI.",
             size=24, bold=True, color=DEEP, line_spacing=1.18)

    box_y = 1.65
    box_h = 4.5

    # Left — donut 51%
    left_x, left_w = 0.55, 5.7
    ocean_box(s, left_x, box_y, left_w, box_h)
    text_box(s, x=left_x + 0.3, y=box_y + 0.2, w=left_w - 0.6, h=0.45,
             text="Проникновение AI в РФ, 2025",
             size=20, bold=True, color=MID, align=PP_ALIGN.CENTER)
    donut_size = 2.65
    donut_x = left_x + (left_w - donut_size) / 2.0
    donut_y = box_y + 0.85
    add_image(s, ASSETS / "charts/c1-vciom-donut.png",
              x=donut_x, y=donut_y, w=donut_size, h=donut_size)
    overlay_h = 0.95
    overlay_y = donut_y + (donut_size - overlay_h) / 2.0 - 0.05
    text_box(s, x=donut_x, y=overlay_y, w=donut_size, h=overlay_h,
             text="51%",
             size=56, bold=True, color=MID, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=left_x + 0.3, y=donut_y + donut_size + 0.05, w=left_w - 0.6, h=0.4,
             text="пользуются AI раз в неделю+",
             size=14, color=DEEP, align=PP_ALIGN.CENTER, bold=True)
    text_box(s, x=left_x + 0.3, y=box_y + box_h - 0.55, w=left_w - 0.6, h=0.4,
             text="ВЦИОМ-Онлайн 13–15 дек 2025  ·  n=3239",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right — bar chart 5 LLMs (NEW: includes Шедеврум, DeepSeek 20% not 43%)
    right_x = left_x + left_w + 0.4
    right_w = 13.333 - 0.55 - right_x
    ocean_box(s, right_x, box_y, right_w, box_h)
    text_box(s, x=right_x + 0.3, y=box_y + 0.2, w=right_w - 0.6, h=0.45,
             text="Использование LLM в РФ — multi-select, 2025",
             size=20, bold=True, color=MID, align=PP_ALIGN.CENTER)
    bar_w = right_w - 0.4
    bar_h = bar_w * 480.0 / 800.0
    if bar_h > box_h - 1.45:
        bar_h = box_h - 1.45
        bar_w = bar_h * 800.0 / 480.0
    bar_x = right_x + (right_w - bar_w) / 2.0
    bar_y = box_y + 0.8
    add_image(s, ASSETS / "charts/c14-llm-shares-rf.png",
              x=bar_x, y=bar_y, w=bar_w, h=bar_h)
    text_box(s, x=right_x + 0.3, y=box_y + box_h - 0.55, w=right_w - 0.6, h=0.4,
             text="*Multi-select: респонденты могли указать несколько. Сумма ≠ 100%.*",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    footer_sources(s, "ВЦИОМ-Онлайн дек 2025 (51%, n=3239)  ·  ВЦИОМ окт 2025 (multi-select shares, n=1600).")
    speaker_notes(s, "Step 2 reveal: «Вы ставили <Х>%, реальность — 51% раз в неделю+». Multi-select disclaimer ОБЯЗАТЕЛЕН на слайде — это share среди использовавших AI, НЕ market share. DeepSeek 43% (Microsoft global downloads) — НЕ показывать на слайде, только в notes как teachable moment про две методологии.")


def build_s05a(p):
    """s05a — instructor card (carry from v3.6)."""
    s = blank(p)
    text_box(s, x=0.55, y=0.55, w=12.3, h=0.9,
             text="Кто я и почему мне это важно.",
             size=28, bold=True, color=DEEP, line_spacing=1.15)
    left_x, card_y, left_w, card_h = 0.55, 1.7, 4.2, 4.65
    ocean_box(s, left_x, card_y, left_w, card_h)
    mono_size = 2.6
    mono_x = left_x + (left_w - mono_size) / 2.0
    mono_y = card_y + 0.55
    add_image(s, ASSETS / "illustrations/monogram-tile.png",
              x=mono_x, y=mono_y, w=mono_size, h=mono_size)
    star_size = 0.4
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(mono_x + mono_size / 2.0 - star_size / 2.0),
                              Inches(card_y + 0.1),
                              Inches(star_size), Inches(star_size))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.fill.background()
    disable_shadow(star)
    text_box(s, x=left_x + 0.3, y=card_y + mono_size + 0.85, w=left_w - 0.6, h=0.5,
             text="[Имя Фамилия]",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=left_x + 0.3, y=card_y + mono_size + 1.4, w=left_w - 0.6, h=0.4,
             text="преподаватель курса",
             size=14, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    items = [
        ("lucide-briefcase-blue.png", "Опыт с AI", "[N лет]",
         "[работа с моделями / проектами]", TEAL),
        ("lucide-lightbulb-blue.png", "Почему этот курс важен", None,
         "[личная мотивация — заполнит преподаватель]", MID),
        ("lucide-users-blue.png", "Что-то о себе", None,
         "[хобби / факт — снижает дистанцию]", MID),
    ]
    right_x = left_x + left_w + 0.5
    right_w = 13.333 - 0.55 - right_x
    item_h = 1.4
    item_gap = 0.18
    for i, (icon, head, prefix_gold, desc, head_color) in enumerate(items):
        y = card_y + i * (item_h + item_gap)
        ocean_box(s, right_x, y, right_w, item_h)
        add_image(s, ASSETS / "icons" / icon,
                  x=right_x + 0.25, y=y + (item_h - 0.7) / 2.0, w=0.7, h=0.7)
        text_box(s, x=right_x + 1.15, y=y + 0.18, w=right_w - 1.3, h=0.5,
                 text=head, size=20, bold=True, color=head_color)
        if prefix_gold is not None:
            text_runs(s, x=right_x + 1.15, y=y + 0.7, w=right_w - 1.3, h=0.6, runs=[
                {"text": prefix_gold, "size": 14, "bold": True, "color": GOLD},
                {"text": "  ", "size": 14, "color": DEEP},
                {"text": desc, "size": 14, "italic": True, "color": LIGHT},
            ], line_spacing=1.3)
        else:
            text_box(s, x=right_x + 1.15, y=y + 0.7, w=right_w - 1.3, h=0.6,
                     text=desc, size=14, italic=True, color=LIGHT, line_spacing=1.3)
    speaker_notes(s, "Шаблон визитки. Преподаватель заполняет 3 пункта о себе перед лекцией. Цель — короткий контакт с аудиторией перед рамкой курса (s05b). Не больше 1 минуты.")


def build_s05b(p):
    """s05b — course frame + central question (carry from v3.6, refs updated)."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.9,
             text="Главный вопрос курса — не «можно ли AI?», а «НУЖНО ли и ГДЕ?».",
             size=24, bold=True, color=DEEP, line_spacing=1.15)
    text_runs(s, x=0.55, y=1.4, w=12.3, h=0.55, runs=[
        {"text": "Стейкс: ", "size": 14, "color": DEEP, "italic": True, "bold": True},
        {"text": "к 2027 — 80% инженерного workforce должно осваивать GenAI ",
         "size": 14, "color": LIGHT, "italic": True},
        {"text": "(Gartner окт 2024)", "size": 14, "color": LIGHT, "italic": True},
        {"text": ".  В РФ — ", "size": 14, "color": LIGHT, "italic": True},
        {"text": "~90% AI-пилотов не доходят до full industrial deployment ",
         "size": 14, "color": LIGHT, "italic": True, "bold": True},
        {"text": "(CNews/Vedomosti/Intellectual Analytics март 2026).",
         "size": 14, "color": LIGHT, "italic": True},
    ])
    funnel_w = 5.9
    funnel_h = funnel_w * 1200.0 / 1500.0
    funnel_x = 0.55
    funnel_y = 2.05
    add_image(s, ASSETS / "diagrams/d2-funnel-v36-clean.png",
              x=funnel_x, y=funnel_y, w=funnel_w, h=funnel_h)
    right_x = funnel_x + funnel_w + 0.4
    right_w = 13.333 - 0.55 - right_x
    box_y = 2.05
    box_h = 4.7
    ocean_box(s, right_x, box_y, right_w, box_h, stroke=TEAL)
    text_runs(s, x=right_x + 0.35, y=box_y + 0.35, w=right_w - 0.7, h=2.6, runs=[
        {"text": "Завтра — ", "size": 24, "color": DEEP, "bold": True},
        {"text": "почти везде", "size": 24, "color": MID, "bold": True},
        {"text": ".\nСегодня — ", "size": 24, "color": DEEP, "bold": True},
        {"text": "почти никто", "size": 24, "color": MID, "bold": True},
        {"text": ".\nКурс — про ", "size": 24, "color": DEEP, "bold": True},
        {"text": "этот разрыв.", "size": 24, "color": GOLD, "bold": True},
    ], line_spacing=1.35)
    filled_rect(s, right_x + 0.35, box_y + 2.95, right_w - 0.7, 0.04, fill=MID)
    text_runs(s, x=right_x + 0.35, y=box_y + 3.15, w=right_w - 0.7, h=1.4, runs=[
        {"text": "Где AI ", "size": 24, "color": DEEP, "bold": True},
        {"text": "работает", "size": 24, "color": MID, "bold": True},
        {"text": ", где —\n", "size": 24, "color": DEEP, "bold": True},
        {"text": "нет", "size": 24, "color": GOLD, "bold": True},
        {"text": ", и как это понять?", "size": 24, "color": DEEP, "bold": True},
    ], line_spacing=1.3)
    speaker_notes(s, "Центральная рамка курса. Возвращаемся в s14, s18, s27. Стейкс обновлён v5: Gartner окт 2024 (80% workforce upskill GenAI by 2027); CNews/Vedomosti/Intellectual Analytics март 2026 (90% не доходят до full industrial deployment). Tone — исследовательский / навигационный.")


# ============================================================
# NEW SLIDES — s06–s29
# ============================================================

def build_s06(p):
    """s06 — Two definitions (academic vs engineering)."""
    s = blank(p)
    slide_title(s, "Академическое определение размыто, инженерное даёт критерий «можно ли сюда AI».")

    col_y, col_h = 1.7, 4.6
    col_w = 5.85
    gap = 0.45

    # Left col — Academic
    left_x = 0.55
    ocean_box(s, left_x, col_y, col_w, col_h)
    add_image(s, ASSETS / "icons/lucide-book-blue.png",
              x=left_x + 0.35, y=col_y + 0.35, w=0.7, h=0.7)
    text_box(s, x=left_x + 1.2, y=col_y + 0.4, w=col_w - 1.4, h=0.5,
             text="АКАДЕМИЧЕСКОЕ", size=14, bold=True, color=MID)
    text_box(s, x=left_x + 0.35, y=col_y + 1.3, w=col_w - 0.7, h=0.6,
             text="Russell & Norvig (2021)",
             size=14, italic=True, color=LIGHT)
    text_runs(s, x=left_x + 0.35, y=col_y + 1.95, w=col_w - 0.7, h=2.5, runs=[
        {"text": "Система, действующая ", "size": 22, "color": DEEP},
        {"text": "рационально", "size": 22, "color": MID, "bold": True},
        {"text": " для достижения целей при имеющихся данных и ограничениях.",
         "size": 22, "color": DEEP},
    ], line_spacing=1.25)

    # Right col — Engineering
    right_x = left_x + col_w + gap
    ocean_box(s, right_x, col_y, col_w, col_h, stroke=TEAL)
    add_image(s, ASSETS / "icons/lucide-wrench-blue.png",
              x=right_x + 0.35, y=col_y + 0.35, w=0.7, h=0.7)
    text_box(s, x=right_x + 1.2, y=col_y + 0.4, w=col_w - 1.4, h=0.5,
             text="ИНЖЕНЕРНОЕ", size=14, bold=True, color=TEAL)
    text_box(s, x=right_x + 0.35, y=col_y + 1.3, w=col_w - 0.7, h=0.6,
             text="рабочее в курсе",
             size=14, italic=True, color=LIGHT)
    text_runs(s, x=right_x + 0.35, y=col_y + 1.95, w=col_w - 0.7, h=2.5, runs=[
        {"text": "Система, принимающая решения ", "size": 22, "color": DEEP},
        {"text": "на основе данных", "size": 22, "color": TEAL, "bold": True},
        {"text": ", а не явных правил «если…то…».",
         "size": 22, "color": DEEP},
    ], line_spacing=1.25)

    # Gold takeaway band at bottom
    gold_callout(s, 0.55, 6.5, 12.25, 0.55,
                 "Если поведение полностью описывается if-else — это не AI в инженерном смысле, даже если маркетинг говорит «умный».",
                 size=15, bold=True)

    speaker_notes(s, "Существует ≥4 канонических определений AI: Turing 1950, McCarthy 1956, Russell&Norvig 2021, ISO/IEC 22989:2022. Все корректны. Используем 2 параллельно: академическое задаёт горизонт, инженерное даёт критерий «задавать ли вопросы про данные, распределение, ошибку обобщения». Спам-фильтр на rule-list — НЕ AI; на классификаторе — AI. Chapter §1.1.")


def build_s07(p):
    """s07 — Timeline with 3 groups + 2017 anchor."""
    s = blank(p)
    # Slightly shorter title to fit one line at 24pt
    slide_title(s, "70 лет AI: три эпохи и точка перелома 2017.", size=26)

    # Timeline base — горизонтальная ось
    tl_y = 4.45
    tl_x_start = 0.7
    tl_x_end = 12.65
    tl_w = tl_x_end - tl_x_start
    filled_rect(s, tl_x_start, tl_y, tl_w, 0.04, fill=LIGHT)

    # 3 group cards (Ocean rounded boxes) above the timeline
    group_h = 2.5
    group_y = 1.65
    group_w = 3.85
    g_gap = 0.15

    # Group 1 — Discovery (1950-80s) — left
    g1_x = 0.55
    ocean_box(s, g1_x, group_y, group_w, group_h)
    text_box(s, x=g1_x + 0.2, y=group_y + 0.1, w=group_w - 0.4, h=0.4,
             text="ОТКРЫТИЯ И ПРАКТИКА", size=12, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(s, x=g1_x + 0.2, y=group_y + 0.4, w=group_w - 0.4, h=0.3,
             text="1950 — 1980-е", size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_runs(s, x=g1_x + 0.2, y=group_y + 0.85, w=group_w - 0.4, h=1.4, runs=[
        {"text": "1950 ", "size": 12, "bold": True, "color": MID},
        {"text": "Тьюринг — Imitation Game", "size": 12, "color": DEEP},
        {"text": "1956 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "Дартмут — McCarthy «AI»", "size": 12, "color": DEEP},
        {"text": "1966 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "Weizenbaum — ELIZA", "size": 12, "color": DEEP},
        {"text": "1980-е ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "экспертные системы", "size": 12, "color": DEEP},
    ], line_spacing=1.35)

    # Group 2 — Winters & breakthroughs — middle
    g2_x = g1_x + group_w + g_gap
    ocean_box(s, g2_x, group_y, group_w, group_h)
    text_box(s, x=g2_x + 0.2, y=group_y + 0.1, w=group_w - 0.4, h=0.4,
             text="ЗИМЫ И ПРОРЫВЫ", size=12, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(s, x=g2_x + 0.2, y=group_y + 0.4, w=group_w - 0.4, h=0.3,
             text="1973 — 2012", size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_runs(s, x=g2_x + 0.2, y=group_y + 0.85, w=group_w - 0.4, h=1.4, runs=[
        {"text": "1974–80 ", "size": 12, "bold": True, "color": MID},
        {"text": "первая зима (Lighthill)", "size": 12, "color": DEEP},
        {"text": "1987–93 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "вторая зима (Lisp Mach.)", "size": 12, "color": DEEP},
        {"text": "1997 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "Deep Blue (200M поз/сек)", "size": 12, "color": DEEP},
        {"text": "2012 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "AlexNet — GPU + DL", "size": 12, "color": DEEP},
    ], line_spacing=1.35)

    # Group 3 — Inflection & explosion — right (with GOLD highlight)
    g3_x = g2_x + group_w + g_gap
    ocean_box(s, g3_x, group_y, group_w, group_h, stroke=GOLD, stroke_pt=2.0)
    text_box(s, x=g3_x + 0.2, y=group_y + 0.1, w=group_w - 0.4, h=0.4,
             text="ПЕРЕЛОМ И ВЗРЫВ", size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, x=g3_x + 0.2, y=group_y + 0.4, w=group_w - 0.4, h=0.3,
             text="2017 — 2026", size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_runs(s, x=g3_x + 0.2, y=group_y + 0.85, w=group_w - 0.4, h=1.4, runs=[
        {"text": "2017  ", "size": 13, "bold": True, "color": GOLD},
        {"text": "«Attention Is All You Need»", "size": 12, "bold": True, "color": DEEP},
        {"text": "Vaswani et al. (160K+ цит.)", "size": 11, "italic": True, "color": LIGHT, "newpara": True},
        {"text": "2022 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "ChatGPT (1M за 5 дней)", "size": 12, "color": DEEP},
        {"text": "2024–26 ", "size": 12, "bold": True, "color": MID, "newpara": True},
        {"text": "reasoning + агенты", "size": 12, "color": DEEP},
    ], line_spacing=1.35)

    # Gold dot on timeline at 2017 position (about 3/4 way)
    dot_x = tl_x_start + tl_w * 0.66
    star_size = 0.45
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(dot_x - star_size/2), Inches(tl_y - 0.18),
                              Inches(star_size), Inches(star_size))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.fill.background()
    disable_shadow(star)
    text_box(s, x=dot_x - 0.5, y=tl_y + 0.18, w=1.0, h=0.3,
             text="2017", size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Year markers on timeline
    for year, frac in [(1950, 0.0), (1980, 0.27), (2012, 0.6), (2022, 0.74), (2026, 1.0)]:
        x = tl_x_start + tl_w * frac
        text_box(s, x=x - 0.4, y=tl_y + 0.18, w=0.8, h=0.3,
                 text=str(year), size=10, color=SLATE, align=PP_ALIGN.CENTER)

    # AI Effect callout at bottom
    gold_callout(s, 0.55, 6.1, 12.25, 0.7,
                 "AI Effect (Tesler): «AI is whatever hasn't been done yet» — как только задача решена, обыватель перестаёт называть её AI.",
                 size=14, bold=True)
    speaker_notes(s, "3 группы — чтобы 21 фактоид не сливался (P2-13). 2017 = настоящая точка перелома (трансформер устранил рекуррентность, дал параллелизм). Авторы Attention основали Cohere/Character.AI/Adept/Sakana. UDIO к авторам Attention отношения не имеет. AI Effect объясняет, почему публичное восприятие AI до 2022 было LLM-центричным.")


def build_s08(p):
    """s08 — 4 axes of AI classification."""
    s = blank(p)
    slide_title(s, "AI-инструмент имеет 4 координаты — задача, модальность, подход, архитектура.")

    # 4 axis cards in 2×2 grid
    grid_y = 1.6
    grid_h = 2.05
    cell_w = 6.0
    cell_gap = 0.25
    row_gap = 0.2

    cells = [
        ("ПО ЗАДАЧЕ", "lucide-target-blue.png",
         "детекция · сегментация · классификация · регрессия · генерация · кластеризация · рекомендация",
         MID),
        ("ПО МОДАЛЬНОСТИ", "lucide-layers-blue.png",
         "число · текст · изображение · видео · аудио · мультимодальные",
         TEAL),
        ("ПО ПОДХОДУ К ОБУЧЕНИЮ", "lucide-brain-blue.png",
         "supervised · unsupervised · self-supervised · RL",
         MID),
        ("ПО АРХИТЕКТУРЕ", "lucide-cpu-blue.png",
         "CNN · RNN/LSTM · Transformer · Diffusion · GNN",
         TEAL),
    ]
    for i, (head, icon, examples, head_color) in enumerate(cells):
        row, col = i // 2, i % 2
        x = 0.55 + col * (cell_w + cell_gap)
        y = grid_y + row * (grid_h + row_gap)
        ocean_box(s, x, y, cell_w, grid_h, stroke=head_color)
        add_image(s, ASSETS / "icons" / icon,
                  x=x + 0.3, y=y + 0.3, w=0.7, h=0.7)
        text_box(s, x=x + 1.2, y=y + 0.35, w=cell_w - 1.4, h=0.5,
                 text=head, size=14, bold=True, color=head_color)
        text_box(s, x=x + 0.35, y=y + 1.15, w=cell_w - 0.7, h=0.85,
                 text=examples, size=14, color=DEEP, italic=True, line_spacing=1.35)

    # Worked example callout — gold
    gold_callout(s, 0.55, 6.05, 12.25, 0.85,
                 "GitHub Copilot:  генерация · текст · self-supervised + RLHF · Transformer-decoder",
                 size=16, bold=True)

    speaker_notes(s, "Карта 4 осей — конкретная модель имеет координаты по всем 4 одновременно. Worked example Copilot прогоняем устно (на слайде показан результат). По задаче — генерация, по модальности — текст (код = последовательность токенов), по подходу — self-supervised pretraining + RLHF, по архитектуре — Transformer decoder-only. Шаблон для любого нового AI. Chapter §1.4.")


def build_s09(p):
    """s09 — Scale of AI in numbers (4 stat tiles + counter-fact)."""
    s = blank(p)
    slide_title(s, "AI стал инфраструктурой за 3 года: 51% разработчиков ежедневно, 46% кода у юзеров Copilot.")

    # 4 stat tiles — slightly smaller to leave room
    tile_y = 1.7
    tile_h = 2.15
    tile_w = 2.95
    tile_gap = 0.18
    tiles = [
        ("900M", "WAU", "ChatGPT февраль 2026", "OpenAI 2026", MID),
        ("51%", "professional daily", "Stack Overflow 2025", "n=49k+, 177 стран", TEAL),
        ("46%", "кода у юзеров Copilot", "GitHub Octoverse 2025", "Java — 61%", MID),
        ("$244–390B", "AI-рынок", "Statista / McKinsey 2025", "разброс — методология", TEAL),
    ]
    for i, (num, label, src1, src2, color) in enumerate(tiles):
        x = 0.55 + i * (tile_w + tile_gap)
        ocean_box(s, x, tile_y, tile_w, tile_h, stroke=color)
        # Smaller font for $244-390B which is wider
        num_size = 36 if "$" in num else 40
        text_box(s, x=x + 0.1, y=tile_y + 0.2, w=tile_w - 0.2, h=0.85,
                 text=num, size=num_size, bold=True, color=color, align=PP_ALIGN.CENTER,
                 line_spacing=1.0)
        text_box(s, x=x + 0.1, y=tile_y + 1.1, w=tile_w - 0.2, h=0.35,
                 text=label, size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.1, y=tile_y + 1.48, w=tile_w - 0.2, h=0.32,
                 text=src1, size=10, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.1, y=tile_y + 1.78, w=tile_w - 0.2, h=0.32,
                 text=src2, size=10, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)

    # Counter-fact band — GOLD
    cf_y = 4.1
    cf_h = 1.0
    filled_rect(s, 0.55, cf_y, 12.25, cf_h, GOLD_TINT, stroke=GOLD, stroke_pt=2.0,
                radius=True, radius_adj=0.08)
    add_image(s, ASSETS / "icons/lucide-alert-triangle-blue.png",
              x=0.85, y=cf_y + 0.15, w=0.7, h=0.7)
    text_runs(s, x=1.7, y=cf_y + 0.13, w=11.0, h=0.45, runs=[
        {"text": "И при этом в РФ ", "size": 18, "color": DEEP, "bold": True},
        {"text": "~90% AI-пилотов не доходят до прода", "size": 18, "color": GOLD, "bold": True},
        {"text": ".", "size": 18, "color": DEEP, "bold": True},
    ])
    text_box(s, x=1.7, y=cf_y + 0.55, w=11.0, h=0.4,
             text="CNews / Vedomosti / Intellectual Analytics, март 2026  ·  30–40% closed без эффекта  ·  7–10% in production",
             size=12, italic=True, color=LIGHT)

    # Trust callout under
    text_runs(s, x=0.55, y=5.35, w=12.25, h=1.2, runs=[
        {"text": "46% разработчиков ", "size": 13, "color": DEEP, "bold": True, "italic": True},
        {"text": "не доверяют точности AI ", "size": 13, "color": MID, "bold": True, "italic": True},
        {"text": "(vs 31% в 2024) — Stack Overflow 2025.",
         "size": 13, "color": LIGHT, "italic": True},
        {"text": "Доверие падает по мере того, как AI становится повседневным.",
         "size": 13, "color": DEEP, "italic": True, "bold": True, "newpara": True},
    ], line_spacing=1.5)

    footer_sources(s, "OpenAI 2026  ·  Stack Overflow Dev Survey 2025  ·  GitHub Octoverse 2025  ·  Statista / McKinsey 2025  ·  CNews/Vedomosti 2026.")
    speaker_notes(s, "AI = инфраструктура. 51% ежедневно, 84% используют/планируют. 46% не доверяют точности — vs 31% год назад. Доверие падает с использованием. РФ — двухслойная картина: пользователи активны, корпоративный слой буксует. 90% не доходят до full industrial deployment. Chapter §2.1.")


def build_s10(p):
    """s10 — DeepSeek moment (timeline with 3 anchor stats)."""
    s = blank(p)
    slide_title(s, "В AI побеждает изобретательность: $5.6M marginal, $589B капотери Nvidia за один день.")

    # Top: trajectory of paradigms
    text_runs(s, x=0.55, y=1.55, w=12.25, h=0.5, runs=[
        {"text": "Траектория: ", "size": 13, "italic": True, "color": LIGHT},
        {"text": "2022 ", "size": 13, "bold": True, "color": MID},
        {"text": "чат → ", "size": 13, "color": DEEP},
        {"text": "2023 ", "size": 13, "bold": True, "color": MID},
        {"text": "зрение → ", "size": 13, "color": DEEP},
        {"text": "2024 ", "size": 13, "bold": True, "color": MID},
        {"text": "рассуждение → ", "size": 13, "color": DEEP},
        {"text": "2025 ", "size": 13, "bold": True, "color": MID},
        {"text": "код → ", "size": 13, "color": DEEP},
        {"text": "2026 ", "size": 13, "bold": True, "color": GOLD},
        {"text": "действие (агенты, MCP)", "size": 13, "color": GOLD, "bold": True},
    ])

    # 3 anchor stats — DeepSeek timeline
    card_y = 2.2
    card_h = 3.4
    card_w = 4.05
    card_gap = 0.18

    anchors = [
        ("26 декабря 2024", "DeepSeek-V3", "$5.6M", "marginal training run",
         "full infra по SemiAnalysis:\n$1.3 — 1.6 млрд", MID),
        ("20 января 2025", "DeepSeek-R1", "97.3%", "MATH-500",
         "открытая reasoning-модель\nуровня OpenAI o1 (96.4%)", TEAL),
        ("27 января 2025", "Nvidia", "−$589B", "капитализации за один день",
         "крупнейшая single-day капотеря\nв истории фондового рынка", GOLD),
    ]
    for i, (date, name, num, label, body, color) in enumerate(anchors):
        x = 0.55 + i * (card_w + card_gap)
        ocean_box(s, x, card_y, card_w, card_h, stroke=color, stroke_pt=2.0)
        text_box(s, x=x + 0.2, y=card_y + 0.2, w=card_w - 0.4, h=0.35,
                 text=date, size=11, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.2, y=card_y + 0.55, w=card_w - 0.4, h=0.4,
                 text=name, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.2, y=card_y + 1.1, w=card_w - 0.4, h=1.0,
                 text=num, size=44, bold=True, color=color, align=PP_ALIGN.CENTER,
                 line_spacing=1.0)
        text_box(s, x=x + 0.2, y=card_y + 2.05, w=card_w - 0.4, h=0.4,
                 text=label, size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.2, y=card_y + 2.55, w=card_w - 0.4, h=0.8,
                 text=body, size=12, color=LIGHT, italic=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Moral callout — gold
    gold_callout(s, 0.55, 5.95, 12.25, 0.85,
                 "Правила меняются быстрее срока обучения инженера. Курс — про устойчивые концепты, не про API.",
                 size=15, bold=True)

    footer_sources(s, "Bloomberg / Reuters 27 января 2025  ·  DeepSeek 2025  ·  SemiAnalysis 2025  ·  Anthropic MCP 2024 (де-факто стандарт).")
    speaker_notes(s, "$5.6M = marginal cost одного training run V3. Full infra $1.3-1.6B (SemiAnalysis). R1 — открытая reasoning-модель уровня o1 (97.3% vs 96.4% MATH-500). Реакция рынка 27 января — Nvidia $589B за день. Мораль: побеждает изобретательность; правила меняются быстрее срока обучения. MCP — Anthropic ноябрь 2024, стандарт 2025-26.")


def build_s11(p):
    """s11 — Layered mental model (model ⊂ chat ⊂ agent ⊂ app)."""
    s = blank(p)
    slide_title(s, "Модель / чат / агент / приложение — это слои, а не альтернативы. Каждый следующий включает предыдущий.")

    # 4 nested rounded boxes (concentric on the left)
    layers_x = 0.55
    layers_y = 1.65
    base_w, base_h = 6.4, 5.2
    # Outer: app (largest)
    ocean_box(s, layers_x, layers_y, base_w, base_h, stroke=TEAL, stroke_pt=2.0)
    text_box(s, x=layers_x + 0.25, y=layers_y + 0.18, w=base_w - 0.5, h=0.4,
             text="ПРИЛОЖЕНИЕ", size=14, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    text_box(s, x=layers_x + 0.25, y=layers_y + 0.55, w=base_w - 0.5, h=0.3,
             text="+ product UX (формы, кнопки, голос)", size=11, italic=True, color=LIGHT)

    # Agent
    inset = 0.45
    a_x = layers_x + inset
    a_y = layers_y + 0.95
    a_w = base_w - 2 * inset
    a_h = base_h - 1.2
    ocean_box(s, a_x, a_y, a_w, a_h, fill=WHITE, stroke=MID, stroke_pt=2.0)
    text_box(s, x=a_x + 0.25, y=a_y + 0.18, w=a_w - 0.5, h=0.4,
             text="АГЕНТ", size=14, bold=True, color=MID)
    text_box(s, x=a_x + 0.25, y=a_y + 0.55, w=a_w - 0.5, h=0.3,
             text="+ инструменты + планирование (act/observe)",
             size=11, italic=True, color=LIGHT)

    # Chat
    c_x = a_x + inset
    c_y = a_y + 0.95
    c_w = a_w - 2 * inset
    c_h = a_h - 1.2
    ocean_box(s, c_x, c_y, c_w, c_h, fill=SURFACE, stroke=LIGHT, stroke_pt=2.0)
    text_box(s, x=c_x + 0.25, y=c_y + 0.18, w=c_w - 0.5, h=0.4,
             text="ЧАТ", size=14, bold=True, color=LIGHT)
    text_box(s, x=c_x + 0.25, y=c_y + 0.55, w=c_w - 0.5, h=0.3,
             text="+ UI + память диалога", size=11, italic=True, color=LIGHT)

    # Model — innermost (gold accent)
    m_x = c_x + inset
    m_y = c_y + 0.95
    m_w = c_w - 2 * inset
    m_h = c_h - 1.2
    ocean_box(s, m_x, m_y, m_w, m_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, x=m_x + 0.2, y=m_y + 0.15, w=m_w - 0.4, h=0.5,
             text="МОДЕЛЬ", size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, x=m_x + 0.2, y=m_y + 0.55, w=m_w - 0.4, h=0.4,
             text="stateless inference,\none task",
             size=11, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.2)

    # Right column — text descriptions (compact)
    right_x = layers_x + base_w + 0.5
    right_w = 13.333 - 0.55 - right_x
    text_runs(s, x=right_x, y=1.65, w=right_w, h=4.55, runs=[
        {"text": "Модель", "size": 17, "bold": True, "color": GOLD},
        {"text": "  =  stateless inference, single-task. Без UI/памяти/инструментов.",
         "size": 13, "color": DEEP},
        {"text": "Чат", "size": 17, "bold": True, "color": LIGHT, "newpara": True},
        {"text": "  =  модель + UI + память диалога. История в контексте.",
         "size": 13, "color": DEEP},
        {"text": "Агент", "size": 17, "bold": True, "color": MID, "newpara": True},
        {"text": "  =  чат + инструменты + планирование. Цикл act/observe.",
         "size": 13, "color": DEEP},
        {"text": "Приложение", "size": 17, "bold": True, "color": TEAL, "newpara": True},
        {"text": "  =  агент или чат + продуктовый UX. AI = функция, не продукт.",
         "size": 13, "color": DEEP},
    ], line_spacing=1.6)

    # Demo announcement
    gold_callout(s, 0.55, 6.55, 12.25, 0.55,
                 "Сейчас увидим один запрос «HTML-страница с графиком» через 3 способа: модель → чат → агент.",
                 size=13, bold=True)

    speaker_notes(s, "Слоистая mental model — главный концепт раздела 3. Не альтернативы — каждый следующий уровень включает предыдущий и добавляет компоненты. Чем выше уровень обвязки, тем больше места для багов между AI-компонентом и его окружением. Анонс s12: одна задача через 3 способа. Chapter §3.1.")


def build_s12(p):
    """s12 — Demo: 3 ways on one task."""
    s = blank(p)
    slide_title(s, "Одна задача, три способа: контроль и стоимость растут с обвязкой.")

    # Task callout in monospace-style
    task_y = 1.55
    task_h = 0.85
    filled_rect(s, 0.55, task_y, 12.25, task_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.08)
    text_runs(s, x=0.85, y=task_y + 0.18, w=11.7, h=0.5, runs=[
        {"text": "ЗАДАЧА:  ", "size": 14, "bold": True, "color": GOLD},
        {"text": "«Создай HTML-страницу с графиком продаж за 2025.»",
         "size": 18, "bold": True, "color": DEEP, "font": FONT_MONO},
    ])

    # 3-column comparison table
    table_y = 2.7
    col_w = 4.05
    col_gap = 0.15
    col_h = 3.6

    cols = [
        ("МОДЕЛЬ", "API (OpenRouter)",
         [("Контроль", "низкий"),
          ("Детерминизм", "высокий"),
          ("Шаги", "1"),
          ("Стоимость", "минимальная"),
          ("Время", "секунды")],
         "сырой JSON / код без структуры", MID),
        ("ЧАТ", "Claude web",
         [("Контроль", "средний"),
          ("Детерминизм", "низкий"),
          ("Шаги", "несколько (диалог)"),
          ("Стоимость", "низкая"),
          ("Время", "минуты")],
         "диалог, уточнение стиля, итерация", LIGHT),
        ("АГЕНТ", "Claude Code",
         [("Контроль", "высокий"),
          ("Детерминизм", "низкий"),
          ("Шаги", "многошаговый цикл"),
          ("Стоимость", "средняя"),
          ("Время", "минуты-час")],
         "читает требования, создаёт файл, правит, проверяет", TEAL),
    ]
    for i, (name, prod, params, summary, color) in enumerate(cols):
        x = 0.55 + i * (col_w + col_gap)
        ocean_box(s, x, table_y, col_w, col_h, stroke=color, stroke_pt=2.0)
        # Header
        text_box(s, x=x + 0.2, y=table_y + 0.18, w=col_w - 0.4, h=0.5,
                 text=name, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.2, y=table_y + 0.62, w=col_w - 0.4, h=0.35,
                 text=prod, size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        # Divider
        filled_rect(s, x + 0.3, table_y + 1.05, col_w - 0.6, 0.02, fill=color)
        # Parameters
        for j, (key, val) in enumerate(params):
            row_y = table_y + 1.15 + j * 0.34
            text_box(s, x=x + 0.25, y=row_y, w=(col_w - 0.5) * 0.45, h=0.3,
                     text=key, size=11, color=LIGHT, bold=False)
            text_box(s, x=x + 0.25 + (col_w - 0.5) * 0.45, y=row_y,
                     w=(col_w - 0.5) * 0.55, h=0.3,
                     text=val, size=11, bold=True, color=DEEP)
        # Summary at bottom
        text_box(s, x=x + 0.2, y=table_y + col_h - 0.55, w=col_w - 0.4, h=0.4,
                 text=summary, size=11, italic=True, color=LIGHT,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    # Bottom: backup note
    text_box(s, x=0.55, y=6.55, w=12.25, h=0.4,
             text="Demo: live + видео-backup. Код: assets/code/three-ways/.",
             size=12, italic=True, color=LIGHT)

    speaker_notes(s, "Hybrid live + video. Зафиксированная задача: «HTML-страница с графиком продаж за 2025». 30 сек screencast model → 30 сек screencast chat → 2 мин live agent (Claude Code). 2 мин разбор таблицы. Чему учит: тот же AI («одна модель внутри») даёт разные результаты в зависимости от обвязки. Выбор обвязки = инженерное решение под задачу. Backup: 3 скриншота + видео.")


def build_s13(p):
    """s13 — Model case: conveyor 10K/hour."""
    s = blank(p)
    slide_title(s, "Поток 10K изображений/час с латентностью 50мс — это standalone-модель, не чат.")

    # Left — case card with metrics
    case_x, case_y, case_w, case_h = 0.55, 1.65, 5.85, 4.3
    ocean_box(s, case_x, case_y, case_w, case_h)
    add_image(s, ASSETS / "icons/lucide-cpu-blue.png",
              x=case_x + 0.3, y=case_y + 0.3, w=0.7, h=0.7)
    text_box(s, x=case_x + 1.2, y=case_y + 0.4, w=case_w - 1.4, h=0.4,
             text="КЕЙС", size=14, bold=True, color=MID)
    text_box(s, x=case_x + 0.3, y=case_y + 1.15, w=case_w - 0.6, h=0.7,
             text="Детектировать дефекты на конвейере",
             size=20, bold=True, color=DEEP, line_spacing=1.25)
    # Metrics (more vertical space, no overlap)
    text_runs(s, x=case_x + 0.3, y=case_y + 1.95, w=case_w - 0.6, h=1.7, runs=[
        {"text": "10 000 ", "size": 22, "bold": True, "color": GOLD},
        {"text": "изображений/час", "size": 13, "color": DEEP},
        {"text": "≤ 50 мс ", "size": 22, "bold": True, "color": GOLD, "newpara": True},
        {"text": "латентность на изделие", "size": 13, "color": DEEP},
        {"text": "edge ", "size": 16, "bold": True, "color": MID, "newpara": True},
        {"text": "развёртывание (без облака)", "size": 13, "color": DEEP},
    ], line_spacing=1.45)
    # Definition pill (separate, lower)
    text_box(s, x=case_x + 0.3, y=case_y + case_h - 0.55, w=case_w - 0.6, h=0.45,
             text="Модель = нейросеть, input→output. Без UI/памяти/инструментов.",
             size=11, italic=True, color=LIGHT, line_spacing=1.3)

    # Right — 4 model examples grid (2×2)
    ex_x, ex_y, ex_w, ex_h = 6.6, 1.65, 6.2, 4.3
    examples = [
        ("YOLO", "детекция объектов", "lucide-target-blue.png"),
        ("Whisper", "распознавание речи", "lucide-message-circle-blue.png"),
        ("Stable Diffusion", "генерация изображений", "lucide-sparkles-blue.png"),
        ("AlphaFold", "3D белка · Нобель 2024", "lucide-flask-conical-blue.png"),
    ]
    cw, ch = 2.95, 2.0
    cgap = 0.2
    for i, (name, role, icon) in enumerate(examples):
        row, col = i // 2, i % 2
        x = ex_x + col * (cw + cgap)
        y = ex_y + row * (ch + cgap)
        color = TEAL if i % 2 == 0 else MID
        ocean_box(s, x, y, cw, ch, stroke=color)
        add_image(s, ASSETS / "icons" / icon,
                  x=x + (cw - 0.7) / 2, y=y + 0.2, w=0.7, h=0.7)
        text_box(s, x=x + 0.1, y=y + 1.05, w=cw - 0.2, h=0.45,
                 text=name, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.1, y=y + 1.5, w=cw - 0.2, h=0.4,
                 text=role, size=11, italic=True, color=LIGHT,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    # When to take / not take — 2 columns at bottom
    btm_y = 6.15
    btm_h = 0.85
    # Left — green-ish (use TEAL for positive)
    filled_rect(s, 0.55, btm_y, 6.05, btm_h, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.0, radius=True, radius_adj=0.1)
    text_runs(s, x=0.75, y=btm_y + 0.13, w=5.7, h=0.6, runs=[
        {"text": "БРАТЬ:  ", "size": 13, "bold": True, "color": TEAL},
        {"text": "высокая нагрузка, одна задача, edge-deployment, встраивание в продукт",
         "size": 12, "color": DEEP, "italic": False},
    ], line_spacing=1.3)
    # Right — gold-ish (warning / not take)
    filled_rect(s, 6.75, btm_y, 6.05, btm_h, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.0, radius=True, radius_adj=0.1)
    text_runs(s, x=6.95, y=btm_y + 0.13, w=5.7, h=0.6, runs=[
        {"text": "НЕ БРАТЬ:  ", "size": 13, "bold": True, "color": GOLD},
        {"text": "разовый запрос, нужен диалог, нет потока однотипных задач",
         "size": 12, "color": DEEP, "italic": False},
    ], line_spacing=1.3)

    speaker_notes(s, "Модель = обученная нейросеть с inference-интерфейсом. Stateless, single-task. Здесь работает инженерное определение AI. Кейс конвейера: 10K/час, 50мс — исключает чат и агент. Если коробка покрывает (Cognex, Keyence, MVTec) — берём её; если нет — строим. Базовый слой layered model. Chapter §3.2.")


def build_s14(p):
    """s14 — Chat case: ТЗ document + LLM shares RF (multi-select)."""
    s = blank(p)
    slide_title(s, "Чат = модель + UI + память. Большинство откатившихся пилотов выбрали чат там, где нужна модель или агент.")

    # Left — case card
    case_x, case_y, case_w, case_h = 0.55, 1.65, 5.6, 4.0
    ocean_box(s, case_x, case_y, case_w, case_h)
    add_image(s, ASSETS / "icons/lucide-message-square-blue.png",
              x=case_x + 0.3, y=case_y + 0.3, w=0.7, h=0.7)
    text_box(s, x=case_x + 1.2, y=case_y + 0.4, w=case_w - 1.4, h=0.4,
             text="КЕЙС", size=14, bold=True, color=MID)
    text_box(s, x=case_x + 0.3, y=case_y + 1.15, w=case_w - 0.6, h=1.4,
             text="Разобрать непонятный нормативный документ и составить чек-лист",
             size=18, bold=True, color=DEEP, line_spacing=1.3)
    text_runs(s, x=case_x + 0.3, y=case_y + 2.7, w=case_w - 0.6, h=1.2, runs=[
        {"text": "Решение:  ", "size": 14, "color": LIGHT, "italic": True},
        {"text": "чат\n", "size": 14, "bold": True, "color": MID},
        {"text": "Не нужна модель ", "size": 12, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "(нет потока). ", "size": 12, "color": DEEP},
        {"text": "Не нужен агент ", "size": 12, "color": LIGHT, "italic": True},
        {"text": "(нет tool use).", "size": 12, "color": DEEP},
    ], line_spacing=1.4)

    # Right — LLM RF bar chart
    chart_x, chart_y, chart_w, chart_h = 6.35, 1.65, 6.45, 4.0
    ocean_box(s, chart_x, chart_y, chart_w, chart_h, stroke=TEAL)
    text_box(s, x=chart_x + 0.2, y=chart_y + 0.15, w=chart_w - 0.4, h=0.4,
             text="LLM в РФ — multi-select",
             size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    bar_w = chart_w - 0.4
    bar_h = bar_w * 480.0 / 800.0
    if bar_h > chart_h - 1.2:
        bar_h = chart_h - 1.2
        bar_w = bar_h * 800.0 / 480.0
    bar_x = chart_x + (chart_w - bar_w) / 2.0
    bar_y = chart_y + 0.6
    add_image(s, ASSETS / "charts/c14-llm-shares-rf.png",
              x=bar_x, y=bar_y, w=bar_w, h=bar_h)
    text_box(s, x=chart_x + 0.2, y=chart_y + chart_h - 0.45, w=chart_w - 0.4, h=0.35,
             text="ВЦИОМ окт 2025, n=1600  ·  можно несколько вариантов",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Gold callback to central question
    gold_callout(s, 0.55, 5.95, 12.25, 0.85,
                 "«Где AI работает» = выбор типа взаимодействия под задачу. Чек-лист в s18.",
                 size=15, bold=True)

    footer_sources(s, "Dam et al. 2024 (LLM-чатботы)  ·  ВЦИОМ окт 2025 multi-select shares (n=1600).")
    speaker_notes(s, "Чат — модель + UI + память диалога. Кейс ТЗ — типовой для чата. Bar chart: ВЦИОМ окт 2025 multi-select. Это share среди использовавших AI, НЕ market share. DeepSeek 20% (НЕ 43%) — 43% это global downloads (Microsoft 2026, telemetry). Большинство откатов = чат там, где нужна модель/агент. Возврат к central question. Chapter §3.3.")


def build_s15(p):
    """s15 — RTC pattern (3 columns: bad / role A / role B)."""
    s = blank(p)
    slide_title(s, "Роль радикально меняет ответ. RTC = Роль + Задача + Контекст — паттерн №1.")

    col_y, col_h = 1.65, 4.4
    col_w = 4.05
    col_gap = 0.15

    cols = [
        ("ПЛОХОЙ ПРОМПТ", LIGHT, GOLD,
         "«Расскажи про AI в медицине»",
         "Расплывчатый текст уровня обзорной статьи."),
        ("РОЛЬ A — McKinsey", MID, MID,
         "«Ты — аналитик McKinsey. Подготовь обзор AI-adoption в российских клиниках за 2025. 5 пунктов с цифрами и оценкой зрелости.»",
         "Структурированный документ с количественными оценками."),
        ("РОЛЬ B — контраст", TEAL, TEAL,
         "«Ты — не-специалист, далёкий от технологий. Тот же вопрос.»",
         "Радикально другой стиль и приоритеты — упрощения, метафоры."),
    ]
    for i, (name, color, accent, prompt, result) in enumerate(cols):
        x = 0.55 + i * (col_w + col_gap)
        ocean_box(s, x, col_y, col_w, col_h, stroke=accent)
        text_box(s, x=x + 0.2, y=col_y + 0.18, w=col_w - 0.4, h=0.4,
                 text=name, size=13, bold=True, color=accent, align=PP_ALIGN.CENTER)
        # Prompt box (pseudo-monospace look — italic)
        text_box(s, x=x + 0.25, y=col_y + 0.7, w=col_w - 0.5, h=2.4,
                 text=prompt, size=13, color=DEEP, italic=True,
                 line_spacing=1.4, font=FONT_BODY)
        # Result label
        text_box(s, x=x + 0.25, y=col_y + col_h - 1.0, w=col_w - 0.5, h=0.3,
                 text="↓ Результат:", size=11, italic=True, color=LIGHT)
        text_box(s, x=x + 0.25, y=col_y + col_h - 0.7, w=col_w - 0.5, h=0.6,
                 text=result, size=12, bold=True, color=accent, line_spacing=1.3)

    # RTC formula bottom — gold large
    formula_y = 6.15
    formula_h = 0.85
    filled_rect(s, 0.55, formula_y, 12.25, formula_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.0, radius=True, radius_adj=0.1)
    text_runs(s, x=0.55, y=formula_y + 0.18, w=12.25, h=0.5, runs=[
        {"text": "RTC  =  ", "size": 22, "bold": True, "color": GOLD},
        {"text": "Роль", "size": 22, "bold": True, "color": MID},
        {"text": "  +  ", "size": 22, "bold": True, "color": GOLD},
        {"text": "Задача", "size": 22, "bold": True, "color": LIGHT},
        {"text": "  +  ", "size": 22, "bold": True, "color": GOLD},
        {"text": "Контекст", "size": 22, "bold": True, "color": TEAL},
    ], align=PP_ALIGN.CENTER)

    speaker_notes(s, "RTC — упрощённый из 16 паттернов White et al. 2023. Роль задаёт регистр; Задача — конкретный продукт; Контекст — данные/ограничения/формат. Контраст 2 ролей — приём отладки. Студенты НЕ пишут — это демо-показ. RTC — паттерн №1 в чит-шите. CoT, Few-Shot, ReAct — позже. Chapter §3.4.")


def build_s16(p):
    """s16 — Agent: 200 PDF case + 5 levels of autonomy."""
    s = blank(p)
    slide_title(s, "Агент = чат + инструменты + планирование. Уровень автономии — ваше дизайн-решение.")

    # Left — case + formula
    left_x, left_w = 0.55, 5.55
    case_y, case_h = 1.65, 2.0
    ocean_box(s, left_x, case_y, left_w, case_h)
    add_image(s, ASSETS / "icons/lucide-file-search-blue.png",
              x=left_x + 0.3, y=case_y + 0.3, w=0.65, h=0.65)
    text_box(s, x=left_x + 1.1, y=case_y + 0.35, w=left_w - 1.3, h=0.4,
             text="КЕЙС", size=13, bold=True, color=MID)
    text_runs(s, x=left_x + 0.3, y=case_y + 1.1, w=left_w - 0.6, h=0.85, runs=[
        {"text": "Прошерстить ", "size": 18, "color": DEEP, "bold": True},
        {"text": "200 PDF", "size": 24, "bold": True, "color": GOLD},
        {"text": " отчётов и собрать таблицу", "size": 18, "color": DEEP, "bold": True},
    ], line_spacing=1.3)

    # Formula box
    fy, fh = case_y + case_h + 0.2, 1.2
    filled_rect(s, left_x, fy, left_w, fh, MID, radius=True, radius_adj=0.12)
    text_box(s, x=left_x + 0.2, y=fy + 0.15, w=left_w - 0.4, h=0.4,
             text="Формула (Weng 2023):",
             size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(s, x=left_x + 0.2, y=fy + 0.5, w=left_w - 0.4, h=0.65,
             text="Agent = LLM + Memory + Planning + Tools",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Examples row
    ex_y = fy + fh + 0.15
    text_runs(s, x=left_x, y=ex_y, w=left_w, h=0.5, runs=[
        {"text": "Примеры: ", "size": 12, "italic": True, "color": LIGHT},
        {"text": "Claude Code · Devin · Operator · AutoGPT · CrewAI · Manus",
         "size": 12, "color": DEEP, "bold": True},
    ])

    # Cycle pattern
    cy_y = ex_y + 0.55
    text_runs(s, x=left_x, y=cy_y, w=left_w, h=0.4, runs=[
        {"text": "Цикл: ", "size": 12, "italic": True, "color": LIGHT},
        {"text": "plan → act → observe → iterate",
         "size": 14, "bold": True, "color": TEAL, "font": FONT_MONO},
    ])

    # Right — 5 levels of autonomy ladder
    right_x = left_x + left_w + 0.4
    right_w = 13.333 - 0.55 - right_x
    levels_y, levels_h = 1.65, 4.7
    ocean_box(s, right_x, levels_y, right_w, levels_h, stroke=TEAL)
    text_box(s, x=right_x + 0.25, y=levels_y + 0.18, w=right_w - 0.5, h=0.4,
             text="5 УРОВНЕЙ АВТОНОМИИ — по роли пользователя",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(s, x=right_x + 0.25, y=levels_y + 0.55, w=right_w - 0.5, h=0.3,
             text="Feng, McDonald & Zhang (2025) arXiv:2506.12469",
             size=10, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # 5 level rows
    levels = [
        ("1", "Operator", "одобряет каждый шаг", "Claude Code «approve each»"),
        ("2", "Collaborator", "работают совместно", "парное прог. с Cursor"),
        ("3", "Consultant", "корректирует план", "Devin предлагает декомпозицию"),
        ("4", "Approver", "утверждает узлы", "агент собирает PR → ваш merge"),
        ("5", "Observer", "только наблюдает итог", "AutoGPT на ночь → отчёт"),
    ]
    row_h = 0.7
    rows_start_y = levels_y + 1.05
    for i, (num, name, role, example) in enumerate(levels):
        ry = rows_start_y + i * row_h
        # Number badge — Gold for level 5 (full autonomy = highlight)
        badge_color = GOLD if num == "5" else MID
        badge = filled_rect(s, right_x + 0.25, ry + 0.05, 0.5, 0.5, badge_color,
                           radius=True, radius_adj=0.5)
        text_box(s, x=right_x + 0.25, y=ry + 0.05, w=0.5, h=0.5,
                 text=num, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=right_x + 0.85, y=ry + 0.02, w=2.0, h=0.35,
                 text=name, size=13, bold=True, color=DEEP)
        text_box(s, x=right_x + 0.85, y=ry + 0.35, w=right_w - 1.1, h=0.32,
                 text=f"{role}  ·  {example}",
                 size=11, italic=True, color=LIGHT, line_spacing=1.2)

    # Gold takeaway at bottom
    gold_callout(s, 0.55, 6.45, 12.25, 0.55,
                 "Уровень автономии = design choice, а не врождённое свойство модели.",
                 size=15, bold=True)

    speaker_notes(s, "Агент = чат + tools + планирование + цикл act/observe. Кейс 200 PDF — естественный для агента (многошаговый, нужны инструменты). 5 уровней Feng et al. 2025 — характеризуются ролью пользователя, не сложностью модели. Один агент в operator vs observer — разные продукты с разными требованиями. Будем возвращаться на семинарах. Chapter §3.5.")


def build_s17(p):
    """s17 — App: Google Translate + 8 logos grid + Copilot ambiguity."""
    s = blank(p)
    slide_title(s, "Приложение = AI, упакованный в продуктовый UX. Большинство уже им пользуются, не зная этого.")

    # Left — case + Translate metric
    left_x, left_w = 0.55, 5.4
    case_y, case_h = 1.65, 2.5
    ocean_box(s, left_x, case_y, left_w, case_h)
    add_image(s, ASSETS / "icons/lucide-globe-blue.png",
              x=left_x + 0.3, y=case_y + 0.3, w=0.65, h=0.65)
    text_box(s, x=left_x + 1.1, y=case_y + 0.35, w=left_w - 1.3, h=0.4,
             text="КЕЙС", size=13, bold=True, color=MID)
    text_box(s, x=left_x + 0.3, y=case_y + 1.1, w=left_w - 0.6, h=0.5,
             text="Перевод техдокументации раз в неделю",
             size=15, bold=True, color=DEEP, line_spacing=1.3)
    text_runs(s, x=left_x + 0.3, y=case_y + 1.7, w=left_w - 0.6, h=0.6, runs=[
        {"text": "Решение:  ", "size": 13, "italic": True, "color": LIGHT},
        {"text": "Google Translate / DeepL", "size": 14, "bold": True, "color": MID},
        {"text": "  (готовое приложение)", "size": 13, "italic": True, "color": LIGHT},
    ])

    # Translate metric — large gold
    metric_y = case_y + case_h + 0.2
    metric_h = 1.7
    filled_rect(s, left_x, metric_y, left_w, metric_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.0, radius=True, radius_adj=0.1)
    text_box(s, x=left_x + 0.2, y=metric_y + 0.15, w=left_w - 0.4, h=0.45,
             text="GOOGLE TRANSLATE, 2026",
             size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, x=left_x + 0.2, y=metric_y + 0.55, w=left_w - 0.4, h=0.7,
             text="1B+ users  ·  ~1T слов/мес",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=left_x + 0.2, y=metric_y + 1.25, w=left_w - 0.4, h=0.4,
             text="across Translate, Search, Lens, Circle to Search",
             size=10, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right — 8 logos grid (3×3 with center "AI inside" tile)
    grid_x, grid_y = 6.2, 1.65
    grid_w, grid_h = 6.6, 4.4
    ocean_box(s, grid_x, grid_y, grid_w, grid_h, stroke=TEAL)
    text_box(s, x=grid_x + 0.2, y=grid_y + 0.18, w=grid_w - 0.4, h=0.4,
             text="ПРИЛОЖЕНИЯ С AI ВНУТРИ",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    apps = [
        "Google Translate", "DeepL", "Grammarly",
        "Notion AI", "Copilot inline", "Яндекс.Навигатор",
        "Adobe Firefly", "Алиса", "Spotify рекомендации",
    ]
    cell_w = (grid_w - 0.6) / 3
    cell_h = (grid_h - 0.95) / 3
    cell_gap = 0.05
    grid_inner_x = grid_x + 0.3
    grid_inner_y = grid_y + 0.7
    for i, app in enumerate(apps):
        row, col = i // 3, i % 3
        x = grid_inner_x + col * cell_w
        y = grid_inner_y + row * cell_h
        # Mini tile
        bg = SURFACE if (row + col) % 2 == 0 else WHITE
        filled_rect(s, x + cell_gap, y + cell_gap,
                    cell_w - 2*cell_gap, cell_h - 2*cell_gap,
                    bg, stroke=LIGHT, stroke_pt=0.5, radius=True, radius_adj=0.15)
        text_box(s, x=x + cell_gap + 0.05, y=y + cell_gap, w=cell_w - 2*cell_gap - 0.1,
                 h=cell_h - 2*cell_gap, text=app, size=11, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    # Copilot ambiguity callout at bottom
    cp_y = 6.25
    cp_h = 0.7
    filled_rect(s, 0.55, cp_y, 12.25, cp_h, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.1)
    text_runs(s, x=0.75, y=cp_y + 0.13, w=11.85, h=0.5, runs=[
        {"text": "GitHub Copilot — особый случай:  ", "size": 14, "bold": True, "color": GOLD},
        {"text": "inline-suggestion (Tab → код) = ", "size": 13, "color": DEEP},
        {"text": "приложение", "size": 13, "bold": True, "color": MID},
        {"text": "  ·  Workspace (план + tool-calls) = ", "size": 13, "color": DEEP},
        {"text": "агент", "size": 13, "bold": True, "color": TEAL},
        {"text": ".", "size": 13, "color": DEEP},
    ])

    speaker_notes(s, "Приложение = AI как функция, не продукт. Детерминированный UX, guardrails, fallbacks. Translate 1T слов/мес — caveat: across Translate/Search/Lens/Circle to Search (P2-fact-2). Копилот: inline = приложение; Workspace = агент. Граница не по бренду, а по архитектуре сценария. Не переплачивайте сложностью. Chapter §3.6.")


def build_s18(p):
    """s18 — Checklist 4 questions + 2x2 quadrant + worked example."""
    s = blank(p)
    slide_title(s, "AI работает там, где задача и тип инструмента совпали. 4 вопроса до внедрения — ваш чек-лист.")

    # Left — 4 questions
    left_x, left_w = 0.55, 6.85
    q_y, q_h = 1.6, 4.5
    ocean_box(s, left_x, q_y, left_w, q_h)
    text_box(s, x=left_x + 0.25, y=q_y + 0.18, w=left_w - 0.5, h=0.4,
             text="ЧЕК-ЛИСТ — задайте до внедрения",
             size=13, bold=True, color=MID)

    questions = [
        ("Q1", "Задача повторяющаяся или разовая?",
         "повторяющаяся → модель / приложение", "разовая → чат", MID),
        ("Q2", "Нужен ли диалог и уточнения?",
         "да → чат / агент", "нет → модель / приложение", LIGHT),
        ("Q3", "Нужны ли внешние инструменты?",
         "да → агент", "нет → модель / чат / приложение", TEAL),
        ("Q4", "Есть ли готовое приложение?",
         "да и покрывает → берём с полки", "нет → строим", GOLD),
    ]
    for i, (qnum, qtext, branch_yes, branch_no, color) in enumerate(questions):
        ry = q_y + 0.7 + i * 0.95
        # Q badge
        badge = filled_rect(s, left_x + 0.25, ry + 0.05, 0.55, 0.55, color,
                           radius=True, radius_adj=0.5)
        text_box(s, x=left_x + 0.25, y=ry + 0.05, w=0.55, h=0.55,
                 text=qnum, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=left_x + 0.95, y=ry + 0.0, w=left_w - 1.2, h=0.4,
                 text=qtext, size=14, bold=True, color=DEEP, line_spacing=1.2)
        text_runs(s, x=left_x + 0.95, y=ry + 0.4, w=left_w - 1.2, h=0.45, runs=[
            {"text": "↳ ", "size": 11, "color": color, "bold": True},
            {"text": branch_yes, "size": 11, "color": DEEP},
            {"text": "    ↳ ", "size": 11, "color": color, "bold": True},
            {"text": branch_no, "size": 11, "color": LIGHT, "italic": True},
        ], line_spacing=1.3)

    # Right — 2×2 matrix
    right_x = left_x + left_w + 0.3
    right_w = 13.333 - 0.55 - right_x
    m_y, m_h = 1.6, 4.5
    ocean_box(s, right_x, m_y, right_w, m_h, stroke=TEAL)
    text_box(s, x=right_x + 0.2, y=m_y + 0.18, w=right_w - 0.4, h=0.4,
             text="МАТРИЦА: Контроль × Детерминированность",
             size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # 2x2 grid inside the box
    inner_pad = 0.5
    grid_w = right_w - 2 * inner_pad
    grid_h = m_h - 1.4
    gx0 = right_x + inner_pad
    gy0 = m_y + 0.85
    cw = grid_w / 2
    ch = grid_h / 2

    # Quadrants (smaller archetype font to avoid wraps; abbreviated label)
    quads = [
        (0, 0, "Высокий контроль\nВысокий детерм.", "МОДЕЛЬ", MID),
        (0, 1, "Высокий UX-контроль\nСредний выход", "ПРИЛО-\nЖЕНИЕ", TEAL),
        (1, 0, "Низкий контроль\nНизкий детерм.", "ЧАТ", LIGHT),
        (1, 1, "Низкий контроль\n+ tools", "АГЕНТ", GOLD),
    ]
    for row, col, label, archetype, color in quads:
        x = gx0 + col * cw
        y = gy0 + row * ch
        ocean_box(s, x + 0.05, y + 0.05, cw - 0.1, ch - 0.1, stroke=color)
        text_box(s, x=x + 0.1, y=y + 0.13, w=cw - 0.2, h=0.6,
                 text=label, size=9, italic=True, color=LIGHT,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)
        text_box(s, x=x + 0.1, y=y + 0.78, w=cw - 0.2, h=0.7,
                 text=archetype, size=14, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

    # Worked example callout — gold
    gold_callout(s, 0.55, 6.2, 12.25, 0.85,
                 "Конвейер 10K/час → Q1 ДА · Q2 НЕТ · Q3 НЕТ · Q4 (если коробки нет) → МОДЕЛЬ.   Полный разбор — методичка §3.8.",
                 size=14, bold=True)

    speaker_notes(s, "Кульминация раздела 3. Порядок Q1-Q4 имеет значение (P1-6): Q1 дешёвый отсев; Q2 широкое ветвление; Q3 узкое; Q4 last-check (не первый). Worked example конвейера: Q1 повторяющаяся → модель/прил.; Q2 нет → исключает чат; Q3 нет → исключает агент; Q4 коробка может покрыть. Возврат к central question. Раздатка ai-choice-checklist.md. Chapter §3.7-§3.8.")


def build_s19(p):
    """s19 — Why boundaries (3 reasons section divider)."""
    s = blank(p)
    slide_title(s, "Границы AI — ваша зона ответственности, не «отдельный вопрос для специалистов по safety».")

    # 3 reason cards — vertical column, but distinct visual
    reasons = [
        ("1", "lucide-shield-blue.png",
         "Вы уже приняли решение встроить AI",
         "Отвечать за инцидент будете не «модель», а вы и ваша команда.",
         MID),
        ("2", "lucide-repeat-blue.png",
         "AI ошибается системно",
         "bias, галлюцинации, sycophancy, distribution shift — не случайные баги, а свойства технологии.",
         TEAL),
        ("3", "lucide-alert-circle-blue.png",
         "Граница «что AI не умеет» — ваша",
         "За пределами — необходимость человеческого решения, верификации, fallback. Не понимаете — система не считается спроектированной.",
         GOLD),
    ]
    card_y = 1.65
    card_h = 1.45
    card_w = 12.25
    card_gap = 0.18
    for i, (num, icon, head, body, color) in enumerate(reasons):
        y = card_y + i * (card_h + card_gap)
        ocean_box(s, 0.55, y, card_w, card_h, stroke=color, stroke_pt=2.0)
        # Number badge
        filled_rect(s, 0.85, y + 0.4, 0.7, 0.7, color, radius=True, radius_adj=0.5)
        text_box(s, x=0.85, y=y + 0.4, w=0.7, h=0.7,
                 text=num, size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Icon
        add_image(s, ASSETS / "icons" / icon,
                  x=1.85, y=y + 0.4, w=0.7, h=0.7)
        # Head
        text_box(s, x=2.85, y=y + 0.25, w=card_w - 2.5, h=0.5,
                 text=head, size=20, bold=True, color=color, line_spacing=1.2)
        # Body
        text_box(s, x=2.85, y=y + 0.8, w=card_w - 2.5, h=0.6,
                 text=body, size=13, italic=True, color=LIGHT, line_spacing=1.35)

    # Bridge phrase — gold
    gold_callout(s, 0.55, 6.7, 12.25, 0.55,
                 "Начнём с того, куда уходят ваши данные.",
                 size=15, bold=True)

    speaker_notes(s, "Раздел 4 = инвентаризация классов проблем, не «глава ужасов». 3 причины — почему ВАША зона: (1) приняли решение, отвечаете; (2) ошибки системны и предсказуемы; (3) границу должны понимать. Структура раздела: 4.2 утечки → 4.3 галлюцинации → 4.4 trio → 4.5 каталог → mid-recap → 4.6 ARC → 4.7 narrow/general → 4.8 Pearl. Chapter §4.1.")


def build_s20(p):
    """s20 — Local vs cloud + Samsung anchor + EU AI Act."""
    s = blank(p)
    slide_title(s, "Consumer-тарифы обучаются на ваших данных по умолчанию. Samsung 2023 — 3 утечки за месяц.")

    # Two columns — Consumer vs Enterprise
    col_y, col_h = 1.6, 3.5
    col_w = 6.0
    col_gap = 0.25

    # Left — Consumer (gold-tinted = warning)
    left_x = 0.55
    ocean_box(s, left_x, col_y, col_w, col_h, stroke=GOLD, stroke_pt=2.0)
    add_image(s, ASSETS / "icons/lucide-cloud-blue.png",
              x=left_x + 0.3, y=col_y + 0.3, w=0.6, h=0.6)
    text_runs(s, x=left_x + 1.05, y=col_y + 0.3, w=col_w - 1.25, h=0.7, runs=[
        {"text": "CONSUMER", "size": 14, "bold": True, "color": GOLD},
        {"text": "  ·  данные → training", "size": 13, "color": GOLD, "italic": True},
    ])
    text_runs(s, x=left_x + 0.3, y=col_y + 1.1, w=col_w - 0.6, h=col_h - 1.3, runs=[
        {"text": "ChatGPT Free / Plus", "size": 13, "bold": True, "color": DEEP},
        {"text": "  → используются для обучения по умолчанию", "size": 12, "color": LIGHT},
        {"text": "Claude consumer", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  → спрашивает с сент. 2025 (5 лет хранения при согласии)",
         "size": 12, "color": LIGHT},
        {"text": "Gemini consumer", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  → данные + ручной ревью до 3 лет (обезличено)",
         "size": 12, "color": LIGHT},
    ], line_spacing=1.55)

    # Right — Enterprise (teal-tinted = safe)
    right_x = left_x + col_w + col_gap
    ocean_box(s, right_x, col_y, col_w, col_h, stroke=TEAL, stroke_pt=2.0)
    add_image(s, ASSETS / "icons/lucide-shield-check-blue.png",
              x=right_x + 0.3, y=col_y + 0.3, w=0.6, h=0.6)
    text_runs(s, x=right_x + 1.05, y=col_y + 0.3, w=col_w - 1.25, h=0.7, runs=[
        {"text": "ENTERPRISE / API", "size": 14, "bold": True, "color": TEAL},
        {"text": "  ·  данные ≠ training", "size": 13, "color": TEAL, "italic": True},
    ])
    text_runs(s, x=right_x + 0.3, y=col_y + 1.1, w=col_w - 0.6, h=col_h - 1.3, runs=[
        {"text": "OpenAI Enterprise / Business / API", "size": 13, "bold": True, "color": DEEP},
        {"text": "  (с марта 2023) → НЕ обучаются", "size": 12, "color": LIGHT},
        {"text": "Anthropic Business", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  → НЕ обучаются", "size": 12, "color": LIGHT},
        {"text": "Google Workspace / Vertex AI", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  → НЕ обучаются  ·  ZDR (Zero Data Retention)",
         "size": 12, "color": LIGHT},
    ], line_spacing=1.55)

    # Samsung incident anchor — gold
    sm_y, sm_h = 5.25, 1.0
    filled_rect(s, 0.55, sm_y, 12.25, sm_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.0, radius=True, radius_adj=0.08)
    add_image(s, ASSETS / "icons/lucide-alert-triangle-blue.png",
              x=0.85, y=sm_y + 0.15, w=0.65, h=0.65)
    text_runs(s, x=1.7, y=sm_y + 0.13, w=11.0, h=0.4, runs=[
        {"text": "SAMSUNG, март–апрель 2023:  ", "size": 14, "bold": True, "color": GOLD},
        {"text": "3 утечки за месяц — код, транскрипт, тестовые последовательности → consumer-ChatGPT.",
         "size": 13, "color": DEEP},
    ])
    text_box(s, x=1.7, y=sm_y + 0.55, w=11.0, h=0.4,
             text="Реакция: запрет внешних AI + лимит 1024 байта на промпт.  Bloomberg, 2023.",
             size=12, italic=True, color=LIGHT)

    # Bottom: alternative + EU AI Act
    btm_y = 6.4
    text_runs(s, x=0.55, y=btm_y, w=12.25, h=0.5, runs=[
        {"text": "Альтернатива: ", "size": 12, "italic": True, "color": LIGHT},
        {"text": "Llama 4 / Mistral / DeepSeek локально через Ollama / LM Studio",
         "size": 12, "color": DEEP, "bold": True},
        {"text": "  ·  breakeven ~100K запросов/день.",
         "size": 12, "italic": True, "color": LIGHT},
    ])
    text_runs(s, x=0.55, y=btm_y + 0.4, w=12.25, h=0.5, runs=[
        {"text": "EU AI Act 2024/1689: ", "size": 12, "italic": True, "color": LIGHT},
        {"text": "стандартный tier до 15M EUR / 3% оборота, верхний (prohibited) — 35M EUR / 7%.",
         "size": 12, "color": DEEP},
    ])

    speaker_notes(s, "Consumer vs Enterprise — критическое различение для безопасности данных. Samsung 2023 — 3 эпизода: код, транскрипт, тестовые последовательности → попали в OpenAI training. Реакция: запрет + лимит 1024 байта. ZDR = Zero Data Retention. Локальное breakeven ~100K req/day. EU AI Act: стандартный tier 15M/3%, prohibited 35M/7% (P2 ревизия). Никогда не загружайте конфиденциальные данные в consumer-AI без проверки политики тарифа. Chapter §4.2.")


def build_s21(p):
    """s21 — Hallucinations: example + Vectara range + anti-pattern."""
    s = blank(p)
    slide_title(s, "AI уверенно генерирует несуществующие DOI. Hallucination rate — от <1% до 15%.")

    # Left — example with prompt + 3 fake DOIs
    left_x, left_w = 0.55, 7.0
    ex_y, ex_h = 1.6, 3.55
    ocean_box(s, left_x, ex_y, left_w, ex_h)
    text_box(s, x=left_x + 0.25, y=ex_y + 0.18, w=left_w - 0.5, h=0.4,
             text="ПРИМЕР", size=13, bold=True, color=MID)
    # Prompt block
    pr_y = ex_y + 0.65
    filled_rect(s, left_x + 0.3, pr_y, left_w - 0.6, 0.8, SURFACE,
                stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.08)
    text_box(s, x=left_x + 0.5, y=pr_y + 0.12, w=left_w - 1.0, h=0.6,
             text="«Назови 3 научные статьи 2023-24 по сейсмостойкости подземных трубопроводов малого диаметра, авторы и DOI.»",
             size=12, italic=True, color=DEEP, line_spacing=1.3, font=FONT_MONO)
    # Result list
    text_runs(s, x=left_x + 0.3, y=pr_y + 0.95, w=left_w - 0.6, h=2.0, runs=[
        {"text": "Ответ AI:\n", "size": 13, "bold": True, "color": LIGHT},
        {"text": "1.  Иванов А. (2023). Журн. инж. геол., DOI: 10.1234/jeg.2023.045  ",
         "size": 12, "color": DEEP, "newpara": True},
        {"text": "❌", "size": 12, "color": GOLD, "bold": True},
        {"text": "2.  Petrov B. (2024). Soil Dynamics, DOI: 10.1016/j.soildyn.108956  ",
         "size": 12, "color": DEEP, "newpara": True},
        {"text": "❌", "size": 12, "color": GOLD, "bold": True},
        {"text": "3.  Smith J., Chen L. (2023). J. Pipeline Sci., DOI: 10.1080/jps.23.0098  ",
         "size": 12, "color": DEEP, "newpara": True},
        {"text": "❌", "size": 12, "color": GOLD, "bold": True},
        {"text": "Журналы — настоящие. Статей нет. DOI не разрешаются.",
         "size": 11, "italic": True, "color": LIGHT, "newpara": True},
    ], line_spacing=1.5)

    # Right — Vectara HHEM bar chart
    right_x = left_x + left_w + 0.3
    right_w = 13.333 - 0.55 - right_x
    h_y, h_h = 1.6, 3.55
    ocean_box(s, right_x, h_y, right_w, h_h, stroke=TEAL)
    text_box(s, x=right_x + 0.2, y=h_y + 0.15, w=right_w - 0.4, h=0.4,
             text="VECTARA HHEM",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(s, x=right_x + 0.2, y=h_y + 0.5, w=right_w - 0.4, h=0.35,
             text="Hallucination rate — диапазон по задачам",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    chart_w = right_w - 0.4
    chart_h = chart_w * 320.0 / 800.0
    if chart_h > h_h - 1.4:
        chart_h = h_h - 1.4
        chart_w = chart_h * 800.0 / 320.0
    add_image(s, ASSETS / "charts/c21-hallucinations.png",
              x=right_x + (right_w - chart_w) / 2, y=h_y + 0.95,
              w=chart_w, h=chart_h)
    text_box(s, x=right_x + 0.2, y=h_y + h_h - 0.45, w=right_w - 0.4, h=0.35,
             text="Vectara 2025-26  ·  github.com/vectara/hallucination-leaderboard",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Anti-pattern — gold large
    ap_y, ap_h = 5.3, 0.95
    filled_rect(s, 0.55, ap_y, 12.25, ap_h, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.1)
    text_runs(s, x=0.85, y=ap_y + 0.13, w=11.55, h=0.5, runs=[
        {"text": "АНТИ-ПАТТЕРН №1:  ", "size": 16, "bold": True, "color": GOLD},
        {"text": "«AI знает всё».", "size": 18, "bold": True, "color": DEEP},
    ])
    text_box(s, x=0.85, y=ap_y + 0.5, w=11.55, h=0.4,
             text="Любой ответ AI по фактическому вопросу = гипотеза, требующая проверки. Особенно — ссылки, цифры, цитаты, нормы.",
             size=13, italic=True, color=DEEP)

    # Sopostav data at bottom
    text_box(s, x=0.55, y=6.45, w=12.25, h=0.45,
             text="~38% сотрудников делятся sensitive info с AI без ведома работодателя.  CybSafe & NCA «Oh Behave!» 2024-25, n=7000, 7 стран.",
             size=12, italic=True, color=LIGHT)

    speaker_notes(s, "Галлюцинация: уверенное порождение фактически неверной информации в форме, неотличимой от верной. Готовый пример: попросить 3 статьи с DOI → проверка через Google Scholar (10-15 сек). Vectara HHEM: <1% (Gemini 2.0 Flash суммаризация) до 10-15% (reasoning). CybSafe ~38% делятся sensitive. Доверие падает — callback к Stack Overflow §2.1. Анти-паттерн «AI знает всё». ПОСЛЕ слайда: retrieval moment s21+ — think-pair-share «найдите подделку», 30 сек в парах. LO7 apply-уровень. Chapter §4.3.")


def build_s22(p):
    """s22 — Bias / Sycophancy / Shift — three concepts + GPT-4o anchor."""
    s = blank(p)
    slide_title(s, "Bias, sycophancy, distribution shift — три проявления одной природы: AI = отражение данных, не источник истины.")

    # 3 cards horizontal
    card_y, card_h = 1.6, 3.45
    card_w = 4.05
    card_gap = 0.15

    cards = [
        ("BIAS", "СМЕЩЕНИЕ",
         "Модель повторяет перекосы датасета.",
         "Пример: HR-скрининг на исторических данных дискриминирует группы, недопредставленные в позитивных примерах.",
         "Самая трудно-исправляемая",
         "lucide-scale-blue.png", MID),
        ("SYCOPHANCY", "ПОДЛИЗЫ",
         "RLHF учит модель соглашаться с пользователем.",
         "Каноник: GPT-4o, апрель 2025. 25 апр release → 28 апр rollback → 29 апр postmortem. «Навязчиво-льстящая».",
         "Самая незаметная",
         "lucide-heart-blue.png", GOLD),
        ("DISTRIBUTION SHIFT", "СДВИГ РАСПРЕДЕЛЕНИЯ",
         "Модель, обученная на 2023, в 2026 уверенно предложит устаревшую библиотеку.",
         "Тихо деградирует со временем без явных сбоев.",
         "Самая частая в долгоживущих системах",
         "lucide-trending-up-blue.png", TEAL),
    ]
    for i, (en_name, ru_name, defn, example, label, icon, color) in enumerate(cards):
        x = 0.55 + i * (card_w + card_gap)
        ocean_box(s, x, card_y, card_w, card_h, stroke=color, stroke_pt=2.0)
        add_image(s, ASSETS / "icons" / icon,
                  x=x + (card_w - 0.7) / 2, y=card_y + 0.18, w=0.7, h=0.7)
        text_box(s, x=x + 0.15, y=card_y + 0.95, w=card_w - 0.3, h=0.4,
                 text=en_name, size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.15, y=card_y + 1.3, w=card_w - 0.3, h=0.4,
                 text=ru_name, size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.15, y=card_y + 1.75, w=card_w - 0.3, h=0.6,
                 text=defn, size=12, color=DEEP, align=PP_ALIGN.LEFT, line_spacing=1.3)
        text_box(s, x=x + 0.15, y=card_y + 2.4, w=card_w - 0.3, h=0.7,
                 text=example, size=11, italic=True, color=LIGHT,
                 line_spacing=1.3)
        # Label band at bottom of card
        text_box(s, x=x + 0.15, y=card_y + card_h - 0.4, w=card_w - 0.3, h=0.3,
                 text=label, size=11, bold=True, color=color,
                 align=PP_ALIGN.CENTER, italic=True)

    # Common cause band — gold
    cc_y, cc_h = 5.2, 1.05
    filled_rect(s, 0.55, cc_y, 12.25, cc_h, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.08)
    text_runs(s, x=0.85, y=cc_y + 0.15, w=11.55, h=0.45, runs=[
        {"text": "ОБЩАЯ ПРИЧИНА:  ", "size": 14, "bold": True, "color": GOLD},
        {"text": "Модель не «знает» истины — воспроизводит закономерности данных.",
         "size": 16, "bold": True, "color": DEEP},
    ])
    text_box(s, x=0.85, y=cc_y + 0.6, w=11.55, h=0.4,
             text="Bias — закономерности истории.  Sycophancy — закономерности RLHF-разметки.  Shift — отсутствие закономерностей нового.",
             size=12, italic=True, color=LIGHT)

    speaker_notes(s, "3 проявления одной природы. Bias — historical perekos (HR-скрининг). Sycophancy — RLHF (определить перед использованием!): люди-разметчики выше оценивают приятные ответы → модель подлизывает. GPT-4o апрель 2025 (P0-8 v5): 25 release, 28 rollback, 29 postmortem. Distribution shift — модель 2023 в 2026 = устаревшие практики. ПОСЛЕ слайда: retrieval moment s22+ — mini-poll «что опаснее в вашей области», LO6 evaluate. Reward hacking / data poisoning / prompt injection — в chapter §4.5 + поздняя лекция safety. Chapter §4.4.")


def build_s23(p):
    """s23 — ARC-AGI economics: 3 bars + open question."""
    s = blank(p)
    slide_title(s, "AI 54% за $30 vs человек 60% за $50–150/час: вопрос — не «хорош ли AI», а «сколько стоит ошибка».")

    # Top — chart in motif box
    chart_x, chart_y = 0.55, 1.55
    chart_w, chart_h = 12.25, 3.0
    ocean_box(s, chart_x, chart_y, chart_w, chart_h, stroke=MID)
    text_box(s, x=chart_x + 0.2, y=chart_y + 0.18, w=chart_w - 0.4, h=0.4,
             text="ARC-AGI-2 — состояние май 2026",
             size=14, bold=True, color=MID, align=PP_ALIGN.CENTER)
    inner_chart_w = chart_w - 0.6
    inner_chart_h = inner_chart_w * 380.0 / 900.0
    if inner_chart_h > chart_h - 0.85:
        inner_chart_h = chart_h - 0.85
        inner_chart_w = inner_chart_h * 900.0 / 380.0
    add_image(s, ASSETS / "charts/c23-arc-agi.png",
              x=chart_x + (chart_w - inner_chart_w) / 2,
              y=chart_y + 0.65, w=inner_chart_w, h=inner_chart_h)

    # Cost annotations under bars
    cost_y = 4.3
    text_runs(s, x=chart_x + 0.55, y=cost_y, w=chart_w - 0.6, h=0.35, runs=[
        {"text": "$50–150/час", "size": 12, "bold": True, "color": DEEP},
        {"text": "                                                ", "size": 12, "color": DEEP},
        {"text": "$30/задачу", "size": 12, "bold": True, "color": GOLD},
        {"text": "                                                ", "size": 12, "color": DEEP},
        {"text": "$2.20/задачу", "size": 12, "bold": True, "color": TEAL},
    ])

    # Open question — gold callout
    oq_y, oq_h = 4.8, 1.2
    filled_rect(s, 0.55, oq_y, 12.25, oq_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.0, radius=True, radius_adj=0.08)
    text_box(s, x=0.85, y=oq_y + 0.15, w=11.55, h=0.4,
             text="ОТКРЫТЫЙ ВОПРОС:",
             size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    text_box(s, x=0.85, y=oq_y + 0.5, w=11.55, h=0.65,
             text="«Если бы вам пришла такая задача — нанять AI с 54% за $30 или человека за $50–150/час? От чего зависит ваш ответ?»",
             size=15, bold=True, italic=True, color=DEEP, line_spacing=1.3)

    # Insight — gold band
    in_y = 6.15
    text_runs(s, x=0.55, y=in_y, w=12.25, h=0.4, runs=[
        {"text": "Инсайт:  ", "size": 14, "italic": True, "color": LIGHT},
        {"text": "Инженерное решение — не «хорош ли AI», а ", "size": 14, "color": DEEP, "bold": True},
        {"text": "«сколько стоит ошибка в моей задаче».",
         "size": 14, "color": GOLD, "bold": True},
    ])
    text_box(s, x=0.55, y=in_y + 0.45, w=12.25, h=0.4,
             text="Источник: arcprize.org (актуальные результаты)  ·  Chollet 2019 (базовая работа).",
             size=11, italic=True, color=LIGHT)

    speaker_notes(s, "ARC-AGI = задачи на абстрактное визуальное рассуждение, простые для людей, трудные для машин. Май 2026: человек 60%, refinement 54% @ $30 (Gemini 3 Pro + Poetiq), single-model 37.6% @ $2.20 (Opus 4.5 Thinking). Открытый вопрос — НЕ gotcha. Ответ «зависит» правильный: цена ошибки, объём, гибридный workflow, регуляторные требования. Инсайт: переключить с «качество модели» на «цена ошибки в моей задаче». Двумерная карта (тип задачи s18 + цена ошибки s23). Chapter §4.6.")


def build_s24(p):
    """s24 — Narrow vs General + 4 leaders predictions spectrum."""
    s = blank(p)
    slide_title(s, "Все существующие AI — narrow. Прогнозы AGI 2-30 лет говорят больше про стимулы, чем про науку.")

    # Top: 2 definitions
    text_runs(s, x=0.55, y=1.55, w=12.25, h=0.7, runs=[
        {"text": "Narrow AI — ", "size": 14, "bold": True, "color": MID},
        {"text": "оптимизирован под конкретную задачу или узкий домен.  ",
         "size": 13, "color": DEEP},
        {"text": "Все существующие production-AI.", "size": 13, "italic": True, "color": LIGHT},
        {"text": "AGI — ", "size": 14, "bold": True, "color": GOLD, "newpara": True},
        {"text": "гипотетический AI с человеческим уровнем когнитивных способностей в широком спектре.  ",
         "size": 13, "color": DEEP},
        {"text": "Открытый вопрос.", "size": 13, "italic": True, "color": LIGHT},
    ], line_spacing=1.5)

    # Spectrum bar with 4 leaders (taller, all labels inside)
    sp_y = 3.05
    sp_h = 3.0
    ocean_box(s, 0.55, sp_y, 12.25, sp_h, stroke=TEAL)
    text_box(s, x=0.75, y=sp_y + 0.15, w=11.85, h=0.4,
             text="ПРОГНОЗЫ AGI — 4 ключевых публичных фигуры",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    # Spectrum line — moved more central
    line_y = sp_y + 1.65
    line_x_start = 1.5
    line_x_end = 11.85
    line_w = line_x_end - line_x_start
    filled_rect(s, line_x_start, line_y, line_w, 0.04, fill=LIGHT)
    # Endpoint labels
    text_box(s, x=line_x_start - 0.6, y=line_y + 0.1, w=1.0, h=0.3,
             text="2 года", size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    text_box(s, x=line_x_end - 0.4, y=line_y + 0.1, w=1.0, h=0.3,
             text="30 лет", size=11, bold=True, color=MID, align=PP_ALIGN.RIGHT)

    # 4 leaders — alternating above/below pattern, all within sp_h
    leaders = [
        (0.05, "S. Altman", "OpenAI", "~5 лет", "AI 100sM users", GOLD),
        (0.30, "D. Amodei", "Anthropic", "2-3 года", "конкурент OpenAI", GOLD),
        (0.65, "D. Hassabis", "Google DeepMind", "50% в декаде", "Нобель 2024, диверсиф.", MID),
        (0.95, "Y. LeCun", "AMI Labs", "не на LLM", "AMI Labs $1B март 2026", TEAL),
    ]
    for i, (frac, name, role, pred, stake, color) in enumerate(leaders):
        x = line_x_start + line_w * frac
        dot_size = 0.25
        filled_rect(s, x - dot_size/2, line_y - dot_size/2 + 0.02,
                   dot_size, dot_size, color, radius=True, radius_adj=0.5)
        # Alternate: even idx (0,2) — above; odd (1,3) — below
        if i % 2 == 0:
            label_y = line_y - 1.0
            stake_y = line_y + 0.25
        else:
            label_y = line_y + 0.3
            stake_y = line_y - 0.4
        text_box(s, x=x - 1.0, y=label_y, w=2.0, h=0.3,
                 text=name, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        text_box(s, x=x - 1.0, y=label_y + 0.3, w=2.0, h=0.25,
                 text=f"({role})", size=9, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        text_box(s, x=x - 1.0, y=label_y + 0.55, w=2.0, h=0.25,
                 text=pred, size=10, color=DEEP, align=PP_ALIGN.CENTER)
        text_box(s, x=x - 1.4, y=stake_y, w=2.8, h=0.25,
                 text=f"stake: {stake}", size=8, italic=True,
                 color=SLATE, align=PP_ALIGN.CENTER)

    # Chinese Room callout — pushed under taller spectrum
    cr_y, cr_h = 6.2, 0.4
    text_runs(s, x=0.85, y=cr_y + 0.05, w=11.55, h=0.35, runs=[
        {"text": "Сёрл (1980), Chinese Room:  ", "size": 11, "bold": True, "color": MID, "italic": True},
        {"text": "симуляция ≠ понимание. Компьютер может выдавать правильные ответы, не понимая их.",
         "size": 11, "italic": True, "color": LIGHT},
    ])

    # Gold takeaway
    gold_callout(s, 0.55, 6.65, 12.25, 0.4,
                 "Вопрос при чтении: «какое решение он бы хотел, чтобы рынок принял?»",
                 size=12, bold=True)

    speaker_notes(s, "Все существующие AI — narrow. AGI открытый вопрос. 4 лидера: Altman ~5 лет (stake: OpenAI valuation); Amodei 2-3 года (stake: конкурент); Hassabis 50% в декаде но «не моделями как сегодня» (Нобель 2024, Google диверсиф.); LeCun «не на LLM» (AMI Labs $1B март 2026, P2-fact-3 ревизия — бывший Meta до ноября 2025). Прогнозы говорят про стимулы. Сёрл 1980. Удалена ссылка на «Давос 2026 LeCun vs Altman» (не верифицирована). Chapter §4.7.")


def build_s25(p):
    """s25 — Human vs AI two columns + Pearl 3 levels pyramid."""
    s = blank(p)
    slide_title(s, "AI работает на уровне 1 (ассоциация); человек думает на 3 (контрфактуальность). Это граница, не починка.")

    # Left — 2 sub-columns (AI better / Human better)
    left_x, left_w = 0.55, 6.5
    col_y, col_h = 1.6, 4.5
    sub_w = (left_w - 0.15) / 2
    # AI better
    ocean_box(s, left_x, col_y, sub_w, col_h, stroke=MID)
    add_image(s, ASSETS / "icons/lucide-cpu-blue.png",
              x=left_x + 0.2, y=col_y + 0.2, w=0.5, h=0.5)
    text_box(s, x=left_x + 0.8, y=col_y + 0.25, w=sub_w - 1.0, h=0.4,
             text="AI ЛУЧШЕ", size=14, bold=True, color=MID)
    text_runs(s, x=left_x + 0.25, y=col_y + 0.85, w=sub_w - 0.5, h=col_h - 1.0, runs=[
        {"text": "•  Скорость", "size": 13, "bold": True, "color": DEEP},
        {"text": "  Deep Blue 200M поз/сек (1997)", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Масштаб", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  AlphaFold 200M структур · Нобель 2024", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Стабильность", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  не устаёт, не отвлекается", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Распознавание паттернов", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  ResNet 3.57% vs human 5.1% ImageNet", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Объём генерации", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  46% кода у юзеров Copilot", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
    ], line_spacing=1.5)

    # Human better
    hb_x = left_x + sub_w + 0.15
    ocean_box(s, hb_x, col_y, sub_w, col_h, stroke=TEAL)
    add_image(s, ASSETS / "icons/lucide-user-blue.png",
              x=hb_x + 0.2, y=col_y + 0.2, w=0.5, h=0.5)
    text_box(s, x=hb_x + 0.8, y=col_y + 0.25, w=sub_w - 1.0, h=0.4,
             text="ЧЕЛОВЕК ЛУЧШЕ", size=14, bold=True, color=TEAL)
    text_runs(s, x=hb_x + 0.25, y=col_y + 0.85, w=sub_w - 0.5, h=col_h - 1.0, runs=[
        {"text": "•  Причинность (Pearl)", "size": 13, "bold": True, "color": DEEP},
        {"text": "  3 уровня — детали справа →", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Абстрактное обобщение", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  по 2-3 примерам (ARC-AGI s23)", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Физический мир", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  парадокс Моравека (1988)", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
        {"text": "•  Целеполагание", "size": 13, "bold": True, "color": DEEP, "newpara": True},
        {"text": "  AI оптимизирует, человек ставит цель", "size": 11, "color": LIGHT, "italic": True, "newpara": True},
    ], line_spacing=1.5)

    # Right — Pearl pyramid
    right_x = 7.35
    right_w = 13.333 - 0.55 - right_x
    py_y, py_h = 1.6, 4.5
    ocean_box(s, right_x, py_y, right_w, py_h, stroke=GOLD, stroke_pt=2.0)
    text_box(s, x=right_x + 0.2, y=py_y + 0.18, w=right_w - 0.4, h=0.4,
             text="PEARL — 3 УРОВНЯ ПРИЧИННОСТИ",
             size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # 3 levels (top→bottom: 3, 2, 1)
    levels_pearl = [
        ("3", "Counterfactual", "что было бы",
         "«был бы наш проект в проде, если бы выбрали другой подход в 2023?»",
         "HUMAN ONLY", GOLD),
        ("2", "Intervention", "что произойдёт",
         "«если поставим лимит $100/мес на API копилота — что с productivity?»",
         "PARTIAL AI", MID),
        ("1", "Association", "что коррелирует",
         "«разработчики с copilot закрывают тикеты на 15% быстрее»",
         "AI ✓", TEAL),
    ]
    lvl_y_start = py_y + 0.7
    lvl_h = (py_h - 0.85) / 3
    for i, (num, en_name, ru_name, example, ai_status, color) in enumerate(levels_pearl):
        ly = lvl_y_start + i * lvl_h
        # Tile
        filled_rect(s, right_x + 0.2, ly, right_w - 0.4, lvl_h - 0.05,
                    SURFACE, stroke=color, stroke_pt=1.5, radius=True, radius_adj=0.1)
        # Number badge
        bd = filled_rect(s, right_x + 0.3, ly + 0.15, 0.45, 0.45, color,
                        radius=True, radius_adj=0.5)
        text_box(s, x=right_x + 0.3, y=ly + 0.15, w=0.45, h=0.45,
                 text=num, size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=right_x + 0.85, y=ly + 0.1, w=right_w - 1.05, h=0.32,
                 text=en_name, size=13, bold=True, color=color)
        text_box(s, x=right_x + 0.85, y=ly + 0.4, w=right_w - 1.05, h=0.3,
                 text=ru_name, size=10, italic=True, color=LIGHT)
        text_box(s, x=right_x + 0.3, y=ly + 0.75, w=right_w - 0.5, h=0.55,
                 text=example, size=10, color=DEEP, italic=True, line_spacing=1.3)
        # Status — small at bottom-right
        text_box(s, x=right_x + 0.3, y=ly + lvl_h - 0.35, w=right_w - 0.5, h=0.25,
                 text=ai_status, size=10, bold=True, color=color, align=PP_ALIGN.RIGHT)

    # Gold takeaway
    gold_callout(s, 0.55, 6.25, 12.25, 0.7,
                 "Решения принимают люди — потому что AI на уровне 1, а решение требует 2-3.",
                 size=15, bold=True)

    speaker_notes(s, "Двухколоночное сопоставление + Pearl пирамида (P1-9 ревизия v2: развёрнутые примеры levels 2 и 3). Уровень 1 — корреляция (LLM умеют). Уровень 2 — что если изменим X (intervention), требует контр. эксперимента или каузальной модели. Уровень 3 — counterfactual «что было бы», LLM не может по принципу — нет каузальной модели реального мира. Pearl объясняет, почему AI-проекты ломаются на стыке аналитики и решения. Решения принимают люди. Chapter §4.8.")


def build_s26(p):
    """s26 — Course roadmap (4 blocks horizontal)."""
    s = blank(p)
    slide_title(s, "17 лекций — 4 блока, путь от концептов к деплою и к границам ответственности.")

    # 4 blocks horizontal
    block_y, block_h = 1.85, 4.4
    block_w = 2.95
    block_gap = 0.18

    blocks = [
        ("БЛОК 1", "Основы", "Лекции 1-4",
         ["Что такое AI", "Трансформер изнутри", "Промпт-инжиниринг", "Базы аналитики"],
         "Классифицируете AI и пишете грамотный промпт", MID),
        ("БЛОК 2", "Инструменты инженера", "Лекции 5-9",
         ["Standalone-модели", "Чат-кейсы", "Агенты", "Приложения с AI", "Patterns"],
         "Выбираете архетип и собираете прототип", LIGHT),
        ("БЛОК 3", "Интеграция в системы", "Лекции 10-13",
         ["MLOps", "Развёртывание", "Мониторинг", "MCP-инфраструктура"],
         "Выводите AI-компонент в продакшн", TEAL),
        ("БЛОК 4", "Границы и ответственность", "Лекции 14-17",
         ["AI safety", "Регуляторика", "Кейсы провалов", "Этика и право"],
         "Проектируете AI-систему с учётом границ", GOLD),
    ]
    for i, (label, name, range_str, topics, outcome, color) in enumerate(blocks):
        x = 0.55 + i * (block_w + block_gap)
        ocean_box(s, x, block_y, block_w, block_h, stroke=color, stroke_pt=2.0)
        # Block label
        text_box(s, x=x + 0.15, y=block_y + 0.15, w=block_w - 0.3, h=0.3,
                 text=label, size=11, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
        # Block name
        text_box(s, x=x + 0.15, y=block_y + 0.45, w=block_w - 0.3, h=0.65,
                 text=name, size=15, bold=True, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)
        # Range
        text_box(s, x=x + 0.15, y=block_y + 1.15, w=block_w - 0.3, h=0.3,
                 text=range_str, size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        # Divider
        filled_rect(s, x + 0.3, block_y + 1.55, block_w - 0.6, 0.02, fill=color)
        # Topics list
        topics_y = block_y + 1.7
        topics_str = "\n".join([f"·  {t}" for t in topics])
        text_box(s, x=x + 0.2, y=topics_y, w=block_w - 0.4, h=2.0,
                 text=topics_str, size=11, color=DEEP, line_spacing=1.5)
        # Outcome at bottom
        outcome_y = block_y + block_h - 0.95
        filled_rect(s, x + 0.15, outcome_y, block_w - 0.3, 0.85,
                    GOLD_TINT if color == GOLD else SURFACE,
                    stroke=color, stroke_pt=0.5, radius=True, radius_adj=0.15)
        text_box(s, x=x + 0.2, y=outcome_y + 0.05, w=block_w - 0.4, h=0.3,
                 text="↳ После блока:", size=9, bold=True, color=color)
        text_box(s, x=x + 0.2, y=outcome_y + 0.3, w=block_w - 0.4, h=0.5,
                 text=outcome, size=10, italic=True, color=DEEP, line_spacing=1.3)

    # Footer note
    text_box(s, x=0.55, y=6.55, w=12.25, h=0.35,
             text="Точные группы и темы — в каталоге `00-course/` (Google Drive); финализируются course-curator до релиза.",
             size=11, italic=True, color=LIGHT)

    speaker_notes(s, "Карта 17 лекций в 4 блока. Outcome-oriented: «после блока — что вы сможете делать». Точные группы и темы вытащить из 00-course/ (Google Drive 1sHXoLaIqCpBRv1IaLjS6lNtBdwI5cPc0) субагентом course-curator до финальной сборки. Зачем roadmap: после 4 разделов лекции 1 студент видит, как фреймворк раскручивается дальше. Без roadmap лекция 1 — оторванный wow-момент.")


def build_s27(p):
    """s27 — Teaser + callback to camera."""
    s = blank(p)
    slide_title(s, "Камера в s1 узнала аудиторию за 30 мс. На лекции 2 — посмотрим КАК.")

    # Left — callback panel
    left_x, left_w = 0.55, 5.6
    cb_y, cb_h = 1.6, 5.0
    ocean_box(s, left_x, cb_y, left_w, cb_h)
    text_box(s, x=left_x + 0.2, y=cb_y + 0.18, w=left_w - 0.4, h=0.4,
             text="CALLBACK", size=13, bold=True, color=MID)

    # Mini YOLO mock
    mini_w = 4.0
    mini_h = mini_w * 720.0 / 1280.0
    add_image(s, ASSETS / "illustrations/s01-yolo-mock.png",
              x=left_x + (left_w - mini_w) / 2, y=cb_y + 0.65, w=mini_w, h=mini_h)
    text_box(s, x=left_x + 0.3, y=cb_y + 0.65 + mini_h + 0.05, w=left_w - 0.6, h=0.3,
             text="s1 — узнала за 30мс / кадр",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Central question recall — gold
    cq_y = cb_y + 0.65 + mini_h + 0.5
    filled_rect(s, left_x + 0.2, cq_y, left_w - 0.4, 1.45, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.1)
    text_box(s, x=left_x + 0.35, y=cq_y + 0.13, w=left_w - 0.7, h=0.35,
             text="Помните вопрос с s5?",
             size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, x=left_x + 0.35, y=cq_y + 0.45, w=left_w - 0.7, h=0.55,
             text="«Где AI работает, где — нет, и как это понять?»",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.25)
    text_box(s, x=left_x + 0.35, y=cq_y + 1.0, w=left_w - 0.7, h=0.4,
             text="Чек-лист s18 = операционализованный ответ.",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right — teaser for L2
    right_x = left_x + left_w + 0.4
    right_w = 13.333 - 0.55 - right_x
    tz_y, tz_h = 1.6, 5.0
    ocean_box(s, right_x, tz_y, right_w, tz_h, stroke=TEAL)
    text_box(s, x=right_x + 0.2, y=tz_y + 0.18, w=right_w - 0.4, h=0.4,
             text="ЛЕКЦИЯ 2 — на уровне чисел",
             size=13, bold=True, color=TEAL)

    # 4 concepts
    concepts = [
        ("lucide-database-blue.png", "Токены", "как AI «видит» текст"),
        ("lucide-network-blue.png", "Эмбеддинги", "числа вместо слов"),
        ("lucide-target-blue.png", "Attention", "что на что смотрит"),
        ("lucide-zap-blue.png", "Температура", "детерминизм vs креативность"),
    ]
    cy_start = tz_y + 0.75
    item_h = 0.85
    for i, (icon, name, desc) in enumerate(concepts):
        cy = cy_start + i * (item_h + 0.1)
        # Mini Ocean box for each item
        filled_rect(s, right_x + 0.2, cy, right_w - 0.4, item_h, SURFACE,
                    stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.18)
        add_image(s, ASSETS / "icons" / icon,
                  x=right_x + 0.35, y=cy + (item_h - 0.55) / 2, w=0.55, h=0.55)
        text_box(s, x=right_x + 1.05, y=cy + 0.1, w=right_w - 1.25, h=0.4,
                 text=name, size=15, bold=True, color=TEAL)
        text_box(s, x=right_x + 1.05, y=cy + 0.45, w=right_w - 1.25, h=0.35,
                 text=desc, size=11, italic=True, color=LIGHT)

    # Gold takeaway
    gold_callout(s, 0.55, 6.7, 12.25, 0.55,
                 "После лекции 2 поймёте, ПОЧЕМУ промпт с ролью работает лучше, и почему AI плохо считает буквы.",
                 size=14, bold=True)

    speaker_notes(s, "Двойной callback. (1) Central question s5: возвращались в s14, s18, сейчас. Чек-лист s18 = операционализованный ответ. (2) Камера s1: на лекции 2 посмотрим КАК — токены/эмбеддинги/attention/температура. Тизер: специфический («ПОЧЕМУ промпт с ролью» + «AI плохо считает буквы»), не общий. Бонус для L2 в notes: Alice in Wonderland как иллюстрация tokenization.")


def build_s28(p):
    """s28 — Three takeaways + homework."""
    s = blank(p)
    slide_title(s, "AI — спектр. Выбор типа — навык. Целеполагание — наше.")

    # 3 takeaway cards horizontal
    card_y, card_h = 1.6, 3.7
    card_w = 4.05
    card_gap = 0.15

    cards = [
        ("1", "lucide-layout-grid-blue.png", "AI — спектр технологий",
         "не одна «умная машина»",
         "4 оси (задача / модальность / подход / архитектура) + 4 архетипа. Грамотное обсуждение начинается с явной классификации.",
         MID),
        ("2", "lucide-check-blue.png", "Выбор типа AI — инженерный навык",
         "и у вас есть инструмент",
         "Чек-лист 4 вопросов из s18 — простой, систематически нарушаемый. Большинство откатов = не задавали эти вопросы.",
         GOLD),
        ("3", "lucide-user-blue.png", "AI усиливает человека",
         "целеполагание, суждение, ответственность — наши",
         "3 категории ошибок (галлюцинации / bias-sycophancy-shift / каталог) требуют человеческого контура. Pearl: AI на уровне 1.",
         TEAL),
    ]
    for i, (num, icon, head, sub, body, color) in enumerate(cards):
        x = 0.55 + i * (card_w + card_gap)
        ocean_box(s, x, card_y, card_w, card_h, stroke=color, stroke_pt=2.0)
        # Number badge
        filled_rect(s, x + (card_w - 0.7) / 2, card_y + 0.2, 0.7, 0.7, color,
                   radius=True, radius_adj=0.5)
        text_box(s, x=x + (card_w - 0.7) / 2, y=card_y + 0.2, w=0.7, h=0.7,
                 text=num, size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Head
        text_box(s, x=x + 0.15, y=card_y + 1.1, w=card_w - 0.3, h=0.85,
                 text=head, size=17, bold=True, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.25)
        # Sub
        text_box(s, x=x + 0.15, y=card_y + 1.95, w=card_w - 0.3, h=0.4,
                 text=sub, size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        # Divider
        filled_rect(s, x + 0.4, card_y + 2.4, card_w - 0.8, 0.02, fill=color)
        # Body
        text_box(s, x=x + 0.2, y=card_y + 2.5, w=card_w - 0.4, h=card_h - 2.6,
                 text=body, size=11, color=DEEP, italic=False, line_spacing=1.4)

    # Homework callout — gold large
    hw_y, hw_h = 5.55, 1.45
    filled_rect(s, 0.55, hw_y, 12.25, hw_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.5, radius=True, radius_adj=0.07)
    add_image(s, ASSETS / "icons/lucide-package-blue.png",
              x=0.85, y=hw_y + 0.2, w=0.7, h=0.7)
    text_box(s, x=1.7, y=hw_y + 0.18, w=11.0, h=0.45,
             text="ДО СЕМИНАРА 1 — домашнее задание",
             size=13, bold=True, color=GOLD)
    text_runs(s, x=1.7, y=hw_y + 0.55, w=11.0, h=0.85, runs=[
        {"text": "Возьмите 1 AI-инструмент, который вы используете регулярно. ",
         "size": 13, "color": DEEP, "bold": True},
        {"text": "Прогоните через 4 вопроса чек-листа s18. Определите архетип. ",
         "size": 13, "color": DEEP},
        {"text": "Оцените 1 потенциальную ошибку из раздела 4. Принесите одностраничный разбор.",
         "size": 13, "color": DEEP, "italic": True},
    ], line_spacing=1.5)

    speaker_notes(s, "3 главных вывода. (1) разделы 1-2; (2) раздел 3 + s18; (3) раздел 4 + Pearl. Ключевая фраза «AI усиливает человека, но целеполагание, критическое суждение и ответственность — наши» — callback к s19 + s25. Домашнее задание apply-уровня: AI-инструмент → 4 вопроса s18 → архетип → 1 ошибка раздела 4 → одностраничный разбор. Покрывает LO1+LO4+LO6+LO7. Tone универсальный.")


def build_s29(p):
    """s29 — Q&A with backup provocations."""
    s = blank(p)
    # No standard title — Q&A is the title

    # Big Q&A in center — gold
    qa_y = 1.0
    qa_h = 2.5
    text_box(s, x=0.55, y=qa_y, w=12.25, h=qa_h,
             text="Q&A",
             size=200, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)

    # Subtitle
    text_box(s, x=0.55, y=qa_y + qa_h + 0.05, w=12.25, h=0.5,
             text="Открытый микрофон.",
             size=22, italic=True, color=DEEP, align=PP_ALIGN.CENTER)

    # 2 backup provocations in motif boxes
    bp_y = 4.6
    bp_h = 1.5
    bp_w = 6.0
    bp_gap = 0.25
    provocations = [
        ("Backup 1",
         "«Поднимите руку, кто через 5 лет не захочет работать с коллегой, который не умеет ставить AI в систему?»",
         MID),
        ("Backup 2",
         "«Кто после этой лекции изменил мнение о чём-то?»",
         TEAL),
    ]
    for i, (label, text, color) in enumerate(provocations):
        x = 0.55 + i * (bp_w + bp_gap)
        ocean_box(s, x, bp_y, bp_w, bp_h, stroke=color)
        text_box(s, x=x + 0.2, y=bp_y + 0.13, w=bp_w - 0.4, h=0.35,
                 text=label, size=12, bold=True, color=color)
        text_box(s, x=x + 0.2, y=bp_y + 0.5, w=bp_w - 0.4, h=0.95,
                 text=text, size=14, italic=True, color=DEEP, line_spacing=1.35)

    # Closing thanks — gold band
    th_y = 6.45
    text_runs(s, x=0.55, y=th_y, w=12.25, h=0.5, runs=[
        {"text": "Спасибо!  ", "size": 18, "bold": True, "color": GOLD},
        {"text": "До семинара — выполните задание s28.",
         "size": 14, "italic": True, "color": LIGHT},
    ], align=PP_ALIGN.CENTER)

    speaker_notes(s, "Открытый Q&A 2 минуты. Если тишина — backup-провокации (универсальный tone, без «ИУ6»): «коллега» вместо «инженер ИУ6». Финал: благодарность + напоминание задание s28. Если Q&A провисает — вернуться к самым обсуждаемым слайдам (s10 DeepSeek / s23 ARC-AGI / s24 narrow vs general).")


# ============================================================
# Main
# ============================================================

def main():
    p = setup_pres()
    builders = [
        build_s01, build_s02, build_s03, build_s04, build_s05a, build_s05b,
        build_s06, build_s07, build_s08, build_s09, build_s10,
        build_s11, build_s12, build_s13, build_s14, build_s15,
        build_s16, build_s17, build_s18,
        build_s19, build_s20, build_s21, build_s22, build_s23, build_s24, build_s25,
        build_s26, build_s27, build_s28, build_s29,
    ]
    for b in builders:
        b(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved: {OUT}  ({len(builders)} slides)")


if __name__ == "__main__":
    main()
