"""
Build script for Семинар 1 — «Знакомство».

Icebreaker seminar, format matched to library/lectures/lec-01 v3 Ocean Gradient
design system (palette, Ocean rounded box motif, typography scale).

Source-of-truth: deck.yaml + slides/*.md.

Canvas: 13.333" x 7.5" (16:9).

=== KNOWN DRIFT (2026-08-08, задокументировано задним числом, NOT YET FIXED) ===

Финальный `rendered/sem-01.pptx` (20 слайдов) был отредактирован владельцем
курса ВРУЧНУЮ поверх обычного pipeline (`deck.yaml`/`slides/*.md` →
`build_sem01.py` → pptx). `deck.yaml` и `slides/*.md` сейчас ТОЧНО описывают
финальный pptx (задача #github-issue: "пересобрать source-файлы задним
числом"), но этот builder-скрипт НЕ обновлён и произведёт СТАРУЮ 19-слайдовую
версию без слайда-биографии, если его запустить как есть. Расхождения:

1. Функции `build_s01`..`build_s19` ниже соответствуют СТАРОЙ нумерации
   (pre-insert). Реальный deck.yaml теперь: s01 (без изменений) → s02
   ИНСТРУКТОР-БИО (НОВЫЙ, нет соответствующей функции) → s03..s20 (сдвиг +1
   относительно старых build_s02..build_s19).
2. Нет функции `build_s02` для нового слайда «О преподавателя» (Левко Максим
   Николаевич) — слайд собран владельцем курса в Google Slides напрямую,
   использует Material Icons ligature-шрифт (send/mail/schema/psychology/
   shield/history_edu), а не Lucide-иконки деки. Полное описание layout —
   в `slides/s02-instructor-bio.md`.
3. `build_s05` (старая функция для stats-слайда, теперь логически s06)
   генерирует старую версию с нижней «живой» 4-слотовой поп-панелью
   («Ваша аудитория») и старыми PNG (`s06-stackoverflow.png`,
   `s06-vciom.png`). Финальный pptx использует `s06-stackoverflow-v2.png` +
   `s06-vciom-v2.png`, НЕ имеет нижней 4-слотовой панели, и содержит
   обновлённую методологию (N=33 244, N=3209, 25–27.06.2026, ±1,7% и т.д.) —
   см. `slides/s06-stats-developers-russia.md`.
4. `build_s02` (старая функция roadmap) использует «Лекции 1–8» / РК на С8;
   финальный pptx (теперь s03) показывает «Лекции 1–6» / «Лекции 7–12» /
   РК1 на С6 — см. `slides/s03-course-roadmap.md`.
5. `build_s03` (старая функция checkpoint mechanics) использует «Семинар 8»;
   финальный pptx (теперь s04) показывает «РК · Семинар 6» — см.
   `slides/s04-checkpoint-mechanics.md`.

TODO (не сделано в этом заходе — приоритет был на source .md/yaml файлах,
см. task brief): переименовать `build_s02`..`build_s19` → `build_s03`..
`build_s20` (сдвиг +1), написать новую `build_s02` под instructor-bio слайд,
переписать содержимое `build_s06` (новая нумерация) под point-fix (убрать
4-слотовую панель, обновить графики/методологию/imagepaths), обновить
`build_s03`/`build_s04` (новая нумерация) под С6-нумерацию РК, обновить
список `BUILDERS`. До этого момента запуск `python3 build_sem01.py`
перезапишет `sem-01.pptx` СТАРОЙ 19-слайдовой версией — НЕ запускать без
предварительного завершения этого рефакторинга.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED Ocean Gradient v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
CODE_BG = RGBColor(0x1B, 0x22, 0x3B)
CODE_FG = RGBColor(0xE3, 0xE9, 0xF2)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons/rendered"
SHOTS = ROOT / "assets/screenshots"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/sem-01.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Courier New"


# ============================================================
# Helpers (adapted from library/lectures/lec-01/rendered/build_lec01.py)
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def multipara_box(slide, x, y, w, h, paragraphs, *,
                   anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """Each item in `paragraphs` is a dict of text_box-style kwargs (text/size/...).
    Uses tf.add_paragraph() per line — the ONLY reliable way to force a line break
    (see notes/mcp-limitations.md [#sem01-render-1]: literal \\n in a single run
    does not line-break reliably under LibreOffice)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, cfg in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", align)
        p.line_spacing = cfg.get("line_spacing", 1.15)
        p.space_after = Pt(cfg.get("space_after", 0))
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", FONT_BODY)
        r.font.size = Pt(cfg.get("size", 14))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                 radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def dashed_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.6,
               radius_pt=12.0, dash="dash"):
    shp = ocean_box(slide, x, y, w, h, fill=fill, stroke=stroke, stroke_pt=stroke_pt,
                     radius_pt=radius_pt)
    ln = shp.line._get_or_add_ln()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    dash_el = etree.SubElement(ln, ns + "prstDash")
    dash_el.set("val", dash)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    path = Path(path)
    if not path.exists():
        print(f"WARNING: missing image {path}")
        return None
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w), height=Inches(h))
    elif w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def icon(slide, name, color_hex, size_px, x, y, w_in):
    """Embed a pre-rendered recolored icon PNG (square) at x,y with width w_in (height = width, square icons)."""
    path = ICONS / f"{name}-{color_hex}-{size_px}.png"
    return add_image(slide, path, x, y, w=w_in, h=w_in)


def slide_title(slide, text, *, y=0.45, h=1.15, w=12.3, x=0.55, size=28,
                color=DEEP, bold=True, line_spacing=1.15, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=14, bold=True):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.06, w=w - 0.4, h=h - 0.12, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.22)


def vote_badge(slide, x, y, *, size=0.42, label=True):
    """Small 'hand + camera' voting-mechanic badge. UNUSED as of round-2 revision —
    the mechanic is explained once on s04 and the badge repeat was judged redundant
    (owner brief round 2). Kept as a helper (not deleted) in case a future slide
    wants it; no call sites remain in BUILDERS as of this revision."""
    icon(slide, "hand", "F0AB00", 64, x, y, size)
    icon(slide, "camera", "21295C", 48, x + size * 0.72, y + size * 0.12, size * 0.62)
    if label:
        text_box(slide, x - 0.1, y + size + 0.03, size + 1.6, 0.3,
                 text="рука + камера", size=9.5, italic=True, color=SLATE,
                 align=PP_ALIGN.LEFT)


def speaker_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r"## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)", md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r"\n+---\s*$", "", notes)
    return notes.strip()


# 7 real Wikimedia Commons photos + 1 best-effort icon tile (cybersecurity — no
# real photo found after honest 6-tier attempt, see iteration-log.md).
INDUSTRY_PHOTO_TILES = [
    ("s01-medicine-davinci.jpg", "Медицина", True),
    ("s01-finance-nyse.jpg", "Финансы", True),
    ("s01-transport-yandex.jpg", "Транспорт", True),
    ("s01-manufacturing-kuka.jpg", "Производство", True),
    ("s01-datacenter.jpg", "Разработка / IT", True),
    ("s01-logistics-ocado.jpg", "Логистика", True),
    ("s01-science-microscope.jpg", "Наука", True),
    ("shield-check", "Кибербезопасность", False),  # icon fallback, best-effort
]


def industry_photo_grid(slide, x, y, w, h, *, highlight_idx=None):
    """8-tile photo collage grid (7 real Wikimedia photos + 1 icon best-effort).
    Used by s01 hero_cover + s19 hero_closing (same asset, same grid — intentional
    callback per brief)."""
    filled_rect(slide, x, y, w, h, DEEP)
    cols, rows = 3, 3
    pad = 0.28
    gap = 0.14
    cell_w = (w - 2 * pad - gap * (cols - 1)) / cols
    cell_h = (h - 2 * pad - gap * (rows - 1)) / rows
    n = len(INDUSTRY_PHOTO_TILES)
    for i, (src, label, is_photo) in enumerate(INDUSTRY_PHOTO_TILES):
        r, c = divmod(i, cols)
        cx = x + pad + c * (cell_w + gap)
        cy = y + pad + r * (cell_h + gap)
        is_hl = (highlight_idx is not None and i == highlight_idx)
        if is_photo:
            path = SHOTS / src
            pic = add_image(slide, path, cx, cy, w=cell_w, h=cell_h - 0.32)
            if pic is not None:
                # crop to fill (avoid distortion): use pptx crop via image aspect mgmt
                pic.left = Inches(cx); pic.top = Inches(cy)
                pic.width = Inches(cell_w); pic.height = Inches(cell_h - 0.32)
            if is_hl:
                gold_frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(cx), Inches(cy), Inches(cell_w), Inches(cell_h - 0.32))
                gold_frame.fill.background()
                gold_frame.line.color.rgb = GOLD
                gold_frame.line.width = Pt(3.0)
                disable_shadow(gold_frame)
        else:
            filled_rect(slide, cx, cy, cell_w, cell_h - 0.32,
                        GOLD if is_hl else MID, radius=True, radius_adj=0.1)
            isz = min(cell_w, cell_h - 0.32) * 0.42
            icon(slide, src, "FFFFFF", 72, cx + (cell_w - isz) / 2,
                 cy + (cell_h - 0.32 - isz) / 2, isz)
        text_box(slide, cx, cy + cell_h - 0.30, cell_w, 0.28,
                 text=label, size=9.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
    text_box(slide, x + pad, y + h - 0.24, w - 2 * pad, 0.2,
             text="Wikimedia Commons · CC / общественное достояние",
             size=8, italic=True, color=RGBColor(0xC8, 0xD2, 0xDF),
             align=PP_ALIGN.RIGHT)


def code_card(slide, x, y, w, h, lines, *, title=None):
    """Dark code-block card (monospace) inside an Ocean rounded box frame."""
    ocean_box(slide, x, y, w, h, fill=SURFACE, stroke=LIGHT)
    pad = 0.2
    inner_y = y + pad
    if title:
        text_box(slide, x + pad, inner_y, w - 2 * pad, 0.35,
                 text=title, size=13, bold=True, color=MID)
        inner_y += 0.42
    code_h = y + h - pad - inner_y
    filled_rect(slide, x + pad, inner_y, w - 2 * pad, code_h, CODE_BG,
                radius=True, radius_adj=0.06)
    paras = [{"text": ln, "size": 12.5, "font": FONT_MONO, "color": CODE_FG,
              "line_spacing": 1.3} for ln in lines]
    multipara_box(slide, x + pad + 0.15, inner_y + 0.14, w - 2 * pad - 0.3,
                  code_h - 0.28, paras)


# ============================================================
# Slide builders
# ============================================================

def build_s01(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    text_box(s, x=0.55, y=0.7, w=6.4, h=0.5,
             text="ДО ПЕРВОЙ ЛЕКЦИИ КУРСА",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.35, w=6.4, h=1.6,
             text="Семинар 1 — знакомство",
             size=44, bold=True, color=DEEP, line_spacing=1.1)
    text_box(s, x=0.55, y=3.15, w=6.4, h=1.0,
             text="Отраслевое применение систем искусственного интеллекта",
             size=17, italic=True, color=MID, line_spacing=1.3)
    text_box(s, x=0.55, y=6.75, w=6.4, h=0.5,
             text="МГТУ им. Н.Э. Баумана", size=12, color=SLATE)
    hero_x, hero_y, hero_w, hero_h = 7.05, 0.0, 6.283, 7.5
    industry_photo_grid(s, hero_x, hero_y, hero_w, hero_h)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "17 лекций × 3 модуля — от основ к разбору AI по отраслям", size=26)

    tl_x, tl_y, tl_w = 0.55, 1.85, 12.23
    tl_h = 1.5
    modules = [
        ("Модуль 1", "Основы и знакомые отрасли", "Лекции 1–8", LIGHT, 8),
        ("Модуль 2", "Высокотехн. производство", "Лекции 9–12", MID, 4),
        ("Модуль 3", "Инфоком, наука, добыча, синтез", "Лекции 13–17", DEEP, 5),
    ]
    total_lec = 17
    mx = tl_x
    for name, desc, rng, col, span in modules:
        mw = tl_w * (span / total_lec)
        filled_rect(s, mx, tl_y, mw, tl_h, col, radius=True, radius_adj=0.08)
        text_box(s, mx + 0.15, tl_y + 0.14, mw - 0.3, 0.35,
                 text=name, size=14, bold=True, color=WHITE)
        text_box(s, mx + 0.15, tl_y + 0.52, mw - 0.3, 0.6,
                 text=desc, size=11, color=WHITE, line_spacing=1.15)
        text_box(s, mx + 0.15, tl_y + tl_h - 0.34, mw - 0.3, 0.3,
                 text=rng, size=10.5, italic=True, color=RGBColor(0xDD, 0xE6, 0xEE))
        mx += mw

    # РК markers at seminar 8/12/17 positions
    checkpoints = [("РК1 · С8", 8), ("РК2 · С12", 12), ("РК3 · С17", 17)]
    for label, pos in checkpoints:
        cx = tl_x + tl_w * (pos / total_lec) - 0.55
        filled_rect(s, cx, tl_y + tl_h + 0.12, 1.1, 0.38, GOLD_TINT, stroke=GOLD,
                    stroke_pt=1.3, radius=True, radius_adj=0.4)
        text_box(s, cx, tl_y + tl_h + 0.12, 1.1, 0.38, text=label, size=10, bold=True,
                 color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # "Сегодня" marker above the timeline start (семинар идёт до Лекции 1 — ни одна
    # лекция не подсвечена как "текущая", только отдельный маркер перед началом)
    chip(s, tl_x, tl_y - 0.5, 1.45, 0.36, "Сегодня", fill=GOLD, color=DEEP, size=12)
    icon(s, "flag", "21295C", 64, tl_x + 1.6, tl_y - 0.47, 0.26)

    # Cross-cutting themes block, visually separated (different fill, own frame)
    # Order (round-2 brief): 1) выбор инструмента 2) типичные ошибки
    # 3) ответственность человек/AI 4) инф. безопасность
    themes_y = tl_y + tl_h + 0.75
    themes_h = 2.35
    ocean_box(s, tl_x, themes_y, tl_w, themes_h, fill=SURFACE, stroke=TEAL)
    text_box(s, tl_x + 0.3, themes_y + 0.18, tl_w - 0.6, 0.4,
             text="Сквозные темы курса — в каждой лекции и семинаре",
             size=15, bold=True, color=DEEP)
    themes = [
        ("compass", "Выбор инструмента", "не любой AI подходит для любой задачи"),
        ("alert-triangle", "Типичные ошибки", "и как их избежать"),
        ("scale", "Ответственность человек / AI", "кто отвечает за результат"),
        ("shield-check", "Инф. безопасность", "что можно и нельзя передавать AI"),
    ]
    tw = (tl_w - 0.6 - 3 * 0.2) / 4
    for i, (ic, title, desc) in enumerate(themes):
        tx = tl_x + 0.3 + i * (tw + 0.2)
        ty = themes_y + 0.75
        icon(s, ic, "028090", 72, tx, ty, 0.42)
        text_box(s, tx, ty + 0.52, tw, 0.4, text=title, size=12, bold=True, color=DEEP,
                 line_spacing=1.1)
        text_box(s, tx, ty + 0.98, tw, 0.55, text=desc, size=10, italic=True,
                 color=SLATE, line_spacing=1.15)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Рубежный контроль — тест дома, разбор на семинаре", size=27)

    # 2-step process with RIGHT_ARROW shapes
    step_y, step_h = 1.9, 2.3
    step_w = 4.6
    gap_arrow = 1.3
    step1_x = 0.9
    step2_x = step1_x + step_w + gap_arrow

    ocean_box(s, step1_x, step_y, step_w, step_h)
    icon(s, "book-open", "065A82", 96, step1_x + 0.3, step_y + 0.3, 0.6)
    text_box(s, step1_x + 0.3, step_y + 1.05, step_w - 0.6, 0.5,
             text="Дома", size=18, bold=True, color=DEEP)
    text_box(s, step1_x + 0.3, step_y + 1.5, step_w - 0.6, 0.7,
             text="Тест на платформе AI-тестирования",
             size=12.5, italic=True, color=SLATE, line_spacing=1.2)

    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
        Inches(step1_x + step_w + 0.15), Inches(step_y + step_h / 2 - 0.28),
        Inches(gap_arrow - 0.3), Inches(0.56))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()
    disable_shadow(arrow)

    ocean_box(s, step2_x, step_y, step_w, step_h)
    icon(s, "users", "028090", 96, step2_x + 0.3, step_y + 0.3, 0.6)
    text_box(s, step2_x + 0.3, step_y + 1.05, step_w - 0.6, 0.5,
             text="На семинаре", size=18, bold=True, color=DEEP)
    text_box(s, step2_x + 0.3, step_y + 1.5, step_w - 0.6, 0.7,
             text="Разбор результатов и открытых вопросов — не сдача теста",
             size=12.5, color=SLATE, line_spacing=1.2)

    # 3 checkpoint chips (same as s02)
    cp_y = step_y + step_h + 0.4
    checkpoints = ["Семинар 8", "Семинар 12", "Семинар 17"]
    cp_w = 2.3
    total_cp_w = cp_w * 3 + 0.3 * 2
    cp_x0 = (SLIDE_W_IN - total_cp_w) / 2
    for i, cptxt in enumerate(checkpoints):
        cx = cp_x0 + i * (cp_w + 0.3)
        filled_rect(s, cx, cp_y, cp_w, 0.55, GOLD_TINT, stroke=GOLD, stroke_pt=1.3,
                    radius=True, radius_adj=0.3)
        text_box(s, cx, cp_y, cp_w, 0.55, text=f"РК · {cptxt}", size=13, bold=True,
                 color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    text_box(s, 0.9, cp_y + 0.85, 11.5, 0.4,
             text="Семинар РК — не сдача теста, а разбор того, что вызвало вопросы",
             size=13, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Поднимите руку на нужный вариант — камера считает", size=25)

    # Explicit mechanic explainer bar
    mech_y = 1.55
    ocean_box(s, 0.55, mech_y, 12.23, 0.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    icon(s, "hand", "F0AB00", 64, 0.85, mech_y + 0.22, 0.5)
    icon(s, "camera", "21295C", 64, 1.45, mech_y + 0.3, 0.38)
    text_box(s, 2.05, mech_y + 0.1, 8.9, 0.5,
             text="Поднимите руку — камера (YOLO) считает и выводит число на экран",
             size=17, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 2.05, mech_y + 0.58, 8.9, 0.32,
             text="не анонимно — число видно всем сразу; отвечаем честно",
             size=11, italic=True, color=SLATE)

    grid_x, grid_y = 0.55, 2.75
    grid_w, grid_h = 12.23, 4.25
    cols, rows = 2, 2
    gap = 0.22
    cw = (grid_w - gap * (cols - 1)) / cols
    ch = (grid_h - gap * (rows - 1)) / rows

    cards = [
        ("sparkles", "Какой AI-инструмент вы используете чаще всего?",
         "ChatGPT · GigaChat/YandexGPT · Copilot/Claude/DeepSeek · не пользуюсь", "4 раунда"),
        ("clock", "Как часто вы используете AI?",
         "каждый день → неделя → месяц → реже → никогда", "5 раундов"),
        ("scale", "Насколько вы доверяете точности AI?",
         "шкала 1–5, раунд на каждое значение", "5 раундов"),
        ("alert-triangle", "Был ли случай, когда AI явно подвёл?",
         "да · нет · не уверен(а)", "3 раунда"),
    ]
    for i, (ic, q, opts, kind) in enumerate(cards):
        r, c = divmod(i, cols)
        cx = grid_x + c * (cw + gap)
        cy = grid_y + r * (ch + gap)
        ocean_box(s, cx, cy, cw, ch)
        pad = 0.22
        icon(s, ic, "065A82", 96, cx + pad, cy + pad, 0.5)
        chip(s, cx + cw - pad - 1.5, cy + pad, 1.5, 0.34, kind, fill=TEAL, size=10)
        text_box(s, cx + pad, cy + pad + 0.62, cw - 2 * pad, 0.6,
                 text=q, size=14.5, bold=True, color=DEEP, line_spacing=1.15)
        text_box(s, cx + pad, cy + ch - pad - 0.6, cw - 2 * pad, 0.6,
                 text=opts, size=11, italic=True, color=SLATE, line_spacing=1.25)
        text_box(s, cx + pad, cy + ch - pad - 0.24, cw - 2 * pad, 0.22,
                 text="раунды поднятия руки, один вариант за раз", size=8.5,
                 italic=True, color=LIGHT)
    speaker_notes(s, load_notes("s04"))


def build_s05(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Используют чаще, доверяют меньше — и это не парадокс", size=23)

    # 2 external-source panels on top (compact) + 1 wide live-audience strip below
    # with a 2x2 grid of the 4 poll questions from s04 — round-2 rebuild: was a
    # 3-panel-in-a-row layout with only 2 live slots, now fits all 4 questions.
    top_y, top_h = 1.5, 2.55
    gap = 0.22
    panel_w = (12.23 - gap) / 2

    # PANEL 1: Stack Overflow
    lx = 0.55
    ocean_box(s, lx, top_y, panel_w, top_h)
    pad = 0.2
    text_box(s, lx + pad, top_y + pad - 0.04, panel_w - 2 * pad, 0.4,
             text="Stack Overflow Developer Survey 2025", size=12.5, bold=True,
             color=MID, line_spacing=1.1)
    chart_path = ASSETS / "charts/s06-stackoverflow.png"
    cw = (panel_w - 2 * pad) * 0.56
    chh = top_h - 2 * pad - 0.38
    chart_y = top_y + pad + 0.34
    if chart_path.exists():
        add_image(s, chart_path, lx + pad, chart_y, w=cw, h=chh)
    text_box(s, lx + pad + cw + 0.15, chart_y, panel_w - 2 * pad - cw - 0.15, chh,
             text="46% явно не доверяют точности (рост с 31%) · 3% доверяют полностью",
             size=9.5, italic=True, color=DEEP, line_spacing=1.25,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + pad, top_y + top_h - pad - 0.24, panel_w - 2 * pad, 0.24,
             text="Источник: Stack Overflow · самоотобранная онлайн-выборка",
             size=8.5, italic=True, color=LIGHT)

    # PANEL 2: VCIOM
    mx = lx + panel_w + gap
    ocean_box(s, mx, top_y, panel_w, top_h)
    text_box(s, mx + pad, top_y + pad - 0.04, panel_w - 2 * pad, 0.4,
             text="ВЦИОМ, «Нейросети в нашей жизни»", size=12.5, bold=True, color=MID,
             line_spacing=1.1)
    chart_path2 = ASSETS / "charts/s06-vciom.png"
    cw2 = (panel_w - 2 * pad) * 0.56
    chh2 = top_h - 2 * pad - 0.38
    chart_y2 = top_y + pad + 0.34
    if chart_path2.exists():
        add_image(s, chart_path2, mx + pad, chart_y2, w=cw2, h=chh2)
    gold_callout(s, mx + pad + cw2 + 0.15, chart_y2, panel_w - 2 * pad - cw2 - 0.15, chh2,
                 "2026: 64% верят в пользу AI · 67% — только в отдельных сферах",
                 size=8.5)
    text_box(s, mx + pad, top_y + top_h - pad - 0.24, panel_w - 2 * pad, 0.24,
             text="Источник: ВЦИОМ · репрезентативная выборка россиян",
             size=8.5, italic=True, color=LIGHT)

    # WIDE LIVE STRIP below: 4 individual mini-cards (own ocean_box outline each)
    # for the 4 poll questions from s04, plain visual placeholders (underscores)
    # only — no "LIVE" tag, no inline instruction text (moved to speaker notes
    # per round-2 brief). Individual cards (not a bare 2x2 grid inside one dashed
    # box) give each question its own visual weight and avoid the excess
    # vertical whitespace a single tall dashed container produced (iter7 review).
    live_y = top_y + top_h + 0.22
    live_h = 7.05 - live_y
    lpad = 0.24
    text_box(s, 0.55, live_y, 4.5, 0.32,
             text="Ваша аудитория", size=13, bold=True, color=DEEP)

    live_qs = [
        ("sparkles", "Инструмент чаще всего"),
        ("clock", "Частота использования"),
        ("scale", "Доверие 1–5"),
        ("alert-triangle", "AI подводил"),
    ]
    cards_y = live_y + 0.4
    cards_h = live_h - 0.4
    gcols = 4
    ggap = 0.18
    gcw = (12.23 - (gcols - 1) * ggap) / gcols
    for i, (ic, lbl) in enumerate(live_qs):
        gx = 0.55 + i * (gcw + ggap)
        dashed_box(s, gx, cards_y, gcw, cards_h, fill=SURFACE, stroke=GOLD, stroke_pt=1.4)
        icon(s, ic, "1C7293", 72, gx + (gcw - 0.44) / 2, cards_y + 0.24, 0.44)
        text_box(s, gx + 0.12, cards_y + 0.84, gcw - 0.24, 0.75,
                 text=lbl, size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.TOP,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)
        text_box(s, gx + 0.12, cards_y + cards_h - 0.55, gcw - 0.24, 0.42,
                 text="____", size=18, bold=True, color=LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "У кого есть история — помогло или подвело?", size=28)

    # Round-2: removed the top "поднимите руку" invite region (icon + text) —
    # the invite mechanic now lives in speaker notes only; freed vertical space
    # goes to a taller category grid instead of empty whitespace.
    grid_x, grid_y = 0.55, 1.75
    grid_w, grid_h = 12.23, 5.15
    cols, rows = 3, 2
    gap = 0.24
    cw = (grid_w - gap * (cols - 1)) / cols
    ch = (grid_h - gap * (rows - 1)) / rows

    cats = [
        ("x-circle", "Галлюцинация факта"),
        ("clock", "Устаревшие данные"),
        ("layers", "Задача решена только\nвнешне — вроде похоже,\nа копнёшь — нет"),
        ("rotate-ccw", "Потеря контекста в диалоге"),
        ("smile", "Неуместный тон"),
        ("target", "Слишком общий совет —\nне под вашу ситуацию"),
    ]
    for i, (ic, lbl) in enumerate(cats):
        r, c = divmod(i, cols)
        cx = grid_x + c * (cw + gap)
        cy = grid_y + r * (ch + gap)
        ocean_box(s, cx, cy, cw, ch)
        icon(s, ic, "028090", 72, cx + (cw - 0.7) / 2, cy + 0.32, 0.7)
        lines = lbl.split("\n")
        paras = [{"text": ln, "size": 13, "bold": True, "color": DEEP,
                  "align": PP_ALIGN.CENTER, "line_spacing": 1.15} for ln in lines]
        multipara_box(s, cx + 0.12, cy + ch - 0.28 - 0.32 * len(lines), cw - 0.24,
                      0.32 * len(lines) + 0.1, paras, align=PP_ALIGN.CENTER)
    # NOTE: no visible "успех-провал-успех" methodology text on slide — CLAUDE.md
    # "No Timing / No Methodology in Slides" rule. That guidance + the "поднимите
    # руку" invite mechanic live in speaker notes only (see load_notes("s06")).
    speaker_notes(s, load_notes("s06"))


# ---- Calibration game (5 categories x 2 slides) ------------------------------
# Round-2 A/B shuffle: fiction + code keep human=A/AI=B (unchanged); technical
# text, image, and dataviz now flip to AI=A/human=B — breaks the "AI is always B"
# predictable pattern (owner brief: minimum 2/5 categories flipped).

def build_s07(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Кто автор — художественный текст: AI или человек?", size=25)

    grid_y = 1.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.9
    texts = [
        "«Как это странно — он умер, а мы живем… Только я подозреваю, что каждый "
        "раз, когда мы ложимся спать, мы точно так же умираем. И солнце уходит "
        "навсегда, и заканчивается вся история. А потом небытие надоедает само "
        "себе, и мы просыпаемся. И мир возникает снова.»",
        "«Сон — это репетиция забвения, которую мы проходим каждую ночь, не "
        "замечая экзамена. Мы гасим свет и на несколько часов перестаём "
        "существовать для самих себя, а утром соглашаемся, что это был перерыв, "
        "а не отмена. Может быть, разница между сном и смертью — только в том, "
        "что после сна нас снова кто-то ждёт.»",
    ]
    for i, txt in enumerate(texts):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.26
        text_box(s, cx + pad, grid_y + pad, cw - 2 * pad, 0.4,
                 text=f"Текст {chr(65+i)}", size=15, bold=True, color=MID)
        # Round-2: removed AI/человек pill chips — voting is by hand-raise
        # (explained on s04), not by clicking a button; text box now uses the
        # freed vertical space (was cut short at ch-2.1 to leave room for chips).
        text_box(s, cx + pad, grid_y + pad + 0.5, cw - 2 * pad, ch - 1.1,
                 text=txt, size=15, italic=True, color=DEEP, line_spacing=1.35,
                 anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Текст A — Пелевин, «Generation П». Текст B — AI", size=24)
    grid_y = 1.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 3.15
    labels = [
        ("book-open", "Виктор Пелевин, «Generation П» (1999)"),
        ("sparkles", "Сгенерировано AI по запросу в похожем стиле"),
    ]
    texts = [
        "«Как это странно — он умер, а мы живем… Только я подозреваю, что каждый "
        "раз, когда мы ложимся спать, мы точно так же умираем. И солнце уходит "
        "навсегда, и заканчивается вся история…»",
        "«Сон — это репетиция забвения, которую мы проходим каждую ночь, не "
        "замечая экзамена. Мы гасим свет и на несколько часов перестаём "
        "существовать для самих себя…»",
    ]
    for i, (txt, (ic, lbl)) in enumerate(zip(texts, labels)):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.26
        icon(s, ic, "028090", 64, cx + pad, grid_y + pad, 0.36)
        text_box(s, cx + pad + 0.48, grid_y + pad - 0.02, cw - 2 * pad - 0.48, 0.5,
                 text=lbl, size=12.5, bold=True, color=MID, line_spacing=1.15,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + pad, grid_y + pad + 0.65, cw - 2 * pad, ch - 1.0,
                 text=txt, size=13.5, italic=True, color=DEEP, line_spacing=1.3)

    # Counter-weight visual for balance: 2 large quote-mark icons framing the conclusion
    concl_y = grid_y + ch + 0.35
    icon(s, "quote", "1C7293", 96, 0.55, concl_y + 0.05, 0.55)
    ocean_box(s, 1.3, concl_y, 10.5, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 1.55, concl_y, 10.0, 0.85,
             text="Стиль можно приблизить, но короткий фрагмент не всегда выдаёт автора",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    icon(s, "quote", "1C7293", 96, 12.0, concl_y + 0.05, 0.55)

    # Extra grounding row: what makes short-fragment style-matching hard vs easy
    tip_y = concl_y + 1.05
    tip_h = 7.2 - tip_y
    if tip_h > 0.6:
        ocean_box(s, 0.55, tip_y, 12.23, tip_h, fill=SURFACE, stroke=TEAL)
        icon(s, "book-open", "028090", 64, 0.85, tip_y + (tip_h - 0.4) / 2, 0.4)
        text_box(s, 1.4, tip_y, 5.3, tip_h,
                 text="Длинное произведение: сюжет, композиция, развитие идеи — сложнее имитировать",
                 size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
        icon(s, "sparkles", "028090", 64, 6.95, tip_y + (tip_h - 0.4) / 2, 0.4)
        text_box(s, 7.5, tip_y, 5.1, tip_h,
                 text="Короткий фрагмент: ритм фразы и интонацию AI имитирует всё убедительнее",
                 size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Кто автор — технический текст: AI или человек?", size=25)

    grid_y = 1.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.9
    # Round-2 A/B shuffle: A=AI, B=human (flipped from round-1 A=human/B=AI).
    texts = [
        "«Конструкция while в Python выполняет тело цикла до тех пор, пока "
        "условие остаётся истинным — в отличие от for, здесь нет встроенного "
        "перебора последовательности: интерпретатор просто проверяет условие "
        "перед каждой итерацией. Это удобно, когда число повторений заранее "
        "неизвестно и зависит от результата вычислений внутри самого цикла.»",
        "«Оператор for используется для перебора элементов последовательности — "
        "например, строки, кортежа или списка — или другого итерируемого объекта. "
        "Список выражений вычисляется один раз и должен вернуть итерируемый "
        "объект; затем набор операторов выполняется один раз для каждого "
        "элемента, в порядке, возвращённом итератором.»",
    ]
    for i, txt in enumerate(texts):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.26
        text_box(s, cx + pad, grid_y + pad, cw - 2 * pad, 0.4,
                 text=f"Текст {chr(65+i)}", size=15, bold=True, color=MID)
        # Round-2: removed AI/человек pill chips (hand-raise voting, not click).
        text_box(s, cx + pad, grid_y + pad + 0.5, cw - 2 * pad, ch - 1.1,
                 text=txt, size=13.5, color=DEEP, line_spacing=1.35,
                 anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Текст A — AI, за секунды. Текст B — официальная документация Python", size=21)
    grid_y = 1.75
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 2.75
    # Round-2 A/B shuffle: A=AI, B=human/official docs (flipped).
    labels = [
        ("sparkles", "AI, по запросу объяснить while в похожем стиле"),
        ("file-code", "Переводы документации Python 3 (digitology.tech)"),
    ]
    texts = [
        "«Конструкция while выполняет тело цикла, пока условие остаётся "
        "истинным — в отличие от for, здесь нет встроенного перебора "
        "последовательности.»",
        "«Оператор for используется для перебора элементов последовательности — "
        "например, строки, кортежа или списка — или другого итерируемого "
        "объекта. Список выражений вычисляется один раз и должен вернуть "
        "итерируемый объект; затем набор операторов выполняется один раз для "
        "каждого элемента, в порядке, возвращённом итератором.»",
    ]
    for i, (txt, (ic, lbl)) in enumerate(zip(texts, labels)):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.24
        icon_color = "F0AB00" if i == 0 else "028090"
        icon(s, ic, icon_color, 64, cx + pad, grid_y + pad, 0.34)
        text_box(s, cx + pad + 0.46, grid_y + pad - 0.02, cw - 2 * pad - 0.46, 0.5,
                 text=lbl, size=11.5, bold=True, color=MID, line_spacing=1.15,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + pad, grid_y + pad + 0.6, cw - 2 * pad, ch - 1.0,
                 text=txt, size=12, color=DEEP, line_spacing=1.25)

    # Tradeoff block — round-2 rewrite: no "not verified, not maintained" AI
    # put-down. New framing: for technical documentation, AI speed + accuracy on
    # standard constructs beats stylistic authorship concerns.
    tw_y = grid_y + ch + 0.3
    tw_h = 1.35
    ocean_box(s, 0.55, tw_y, 12.23, tw_h, fill=SURFACE, stroke=TEAL)
    half = (12.23 - 0.6) / 2
    icon(s, "sparkles", "F0AB00", 64, 0.55 + 0.3, tw_y + (tw_h - 0.34) / 2, 0.34)
    text_box(s, 0.55 + 0.75, tw_y + 0.16, half - 0.45, 0.32,
             text="AI-версия", size=13.5, bold=True, color=GOLD)
    text_box(s, 0.55 + 0.75, tw_y + 0.5, half - 0.45, tw_h - 0.62,
             text="Быстрее и часто не менее точна для стандартных языковых конструкций · экономит время поиска",
             size=12, color=DEEP, line_spacing=1.25)
    icon(s, "file-code", "065A82", 64, 0.55 + 0.3 + half, tw_y + (tw_h - 0.34) / 2, 0.34)
    text_box(s, 0.55 + 0.75 + half, tw_y + 0.16, half - 0.45, 0.32,
             text="Официальная документация", size=13.5, bold=True, color=MID)
    text_box(s, 0.55 + 0.75 + half, tw_y + 0.5, half - 0.45, tw_h - 0.62,
             text="Пишется и поддерживается командой · стабильна годами",
             size=12, color=DEEP, line_spacing=1.25)

    # Counter-weight grounding row (mirrors s08/s12/s16 pattern) — closes remaining
    # whitespace below the tradeoff block and reinforces the lesson visually.
    tip_y = tw_y + tw_h + 0.3
    tip_h = 7.2 - tip_y
    if tip_h > 0.7:
        ocean_box(s, 0.55, tip_y, 12.23, tip_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
        text_box(s, 0.85, tip_y, 11.6, tip_h,
                 text="Урок: для технической документации важна не стилистика и не то, кто автор, "
                      "а точность и скорость ответа — здесь AI особенно силён на типовых вопросах",
                 size=14.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Кто автор — код: AI или человек?", size=27)

    grid_y = 1.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.9
    code_a = [
        "def prepare_method(self, method):",
        '    """Prepares the given HTTP method."""',
        "    self.method = method",
        "    if self.method is not None:",
        "        self.method = to_native_string(",
        "            self.method.upper())",
    ]
    code_b = [
        "def normalize_header_name(name):",
        '    """Normalize an HTTP header name."""',
        "    if not name:",
        "        return name",
        "    return '-'.join(",
        "        part.capitalize() for part in name.split('-'))",
    ]
    # Round-2: removed AI/человек pill chips below each code card (hand-raise
    # voting, not click) — code cards now use the full card height.
    for i, code in enumerate([code_a, code_b]):
        cx = 0.55 + i * (cw + gap)
        code_card(s, cx, grid_y, cw, ch, code, title=f"Код {chr(65+i)}")
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Код A — библиотека requests, 10+ лет в проде. Код B — AI, 10-20 секунд", size=20)
    grid_y = 1.7
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 3.3
    code_a = [
        "def prepare_method(self, method):",
        '    """Prepares the given HTTP method."""',
        "    self.method = method",
        "    if self.method is not None:",
        "        self.method = to_native_string(",
        "            self.method.upper())",
    ]
    code_b = [
        "def normalize_header_name(name):",
        '    """Normalize an HTTP header name."""',
        "    if not name:",
        "        return name",
        "    return '-'.join(",
        "        part.capitalize() for part in name.split('-'))",
    ]
    labels = ["requests, github.com/psf/requests, MIT license",
              "AI, по запросу написать похожую функцию"]
    for i, (code, lbl) in enumerate(zip([code_a, code_b], labels)):
        cx = 0.55 + i * (cw + gap)
        code_card(s, cx, grid_y, cw, ch, code, title=lbl)

    # Tradeoff block — round-2 rewrite: no AI put-down ("without operational
    # history or guarantee of edge cases"). New framing: typical/boilerplate
    # tasks -> AI wins on speed; complex/high-stakes/unique logic -> human.
    tw_y = grid_y + ch + 0.25
    tw_h = 7.2 - tw_y
    ocean_box(s, 0.55, tw_y, 12.23, tw_h, fill=SURFACE, stroke=TEAL)
    half = (12.23 - 0.6) / 2
    text_box(s, 0.85, tw_y + 0.16, half, 0.35, text="requests", size=13.5, bold=True, color=MID)
    text_box(s, 0.85, tw_y + 0.55, half, tw_h - 0.85,
             text="Используется в миллионах проектов · проверена на реальной нагрузке 10+ лет",
             size=12, color=DEEP, line_spacing=1.25)
    text_box(s, 0.85 + half, tw_y + 0.16, half, 0.35, text="AI-версия", size=13.5,
             bold=True, color=GOLD)
    text_box(s, 0.85 + half, tw_y + 0.55, half, tw_h - 0.85,
             text="Отлично подходит для типовых вспомогательных функций за секунды",
             size=12, color=DEEP, line_spacing=1.25)
    gold_callout(s, 0.85, tw_y + tw_h - 0.52, 12.23 - 0.6, 0.42,
                 "Урок: типовые задачи — AI, сложные и уникальные решения, от которых "
                 "зависит стабильность системы, — человеку", size=11.5)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Кто автор — картина: AI или человек?", size=27)

    grid_y = 1.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.9

    # Round-2 A/B shuffle: A=AI (stylized landscape mockup image), B=real
    # painting (flipped from round-1 A=painting/B=AI-text). Round-2 also
    # removed AI/человек pill chips below each card (hand-raise voting, not
    # click) — cards use full height.
    # Round-3 point fix: card A was bare icon+text while card B was a
    # full-bleed image — the format asymmetry alone gave away the answer
    # before voting (student-simulator P1). Fixed by rendering card A as a
    # full-bleed stylized landscape mockup image, same footprint as card B,
    # so both cards are visually indistinguishable as "a picture" at a glance.
    cx = 0.55
    ocean_box(s, cx, grid_y, cw, ch)
    pad = 0.22
    text_box(s, cx + pad, grid_y + pad - 0.02, cw - 2 * pad, 0.35,
             text="Изображение A", size=14, bold=True, color=MID)
    img_path_a = ASSETS / "illustrations" / "s13-ai-landscape-mockup.png"
    img_h = ch - 0.6
    add_image(s, img_path_a, cx + pad, grid_y + pad + 0.4, w=cw - 2 * pad, h=img_h)

    # Image B: real Velasco painting
    cx2 = cx + cw + gap
    ocean_box(s, cx2, grid_y, cw, ch)
    text_box(s, cx2 + pad, grid_y + pad - 0.02, cw - 2 * pad, 0.35,
             text="Изображение B", size=14, bold=True, color=MID)
    img_path = SHOTS / "s13-velasco-real.jpg"
    add_image(s, img_path, cx2 + pad, grid_y + pad + 0.4, w=cw - 2 * pad, h=img_h)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Изображение B — Хосе Мария Веласко, 1875. AI сгенерирует похожее за секунды", size=18)
    grid_y = 1.7
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 5.2

    # Round-2 A/B shuffle: A=AI (text card), B=real painting.
    cx = 0.55
    ocean_box(s, cx, grid_y, cw, ch, fill=SURFACE, stroke=TEAL)
    icon(s, "sparkles", "1C7293", 96, cx + (cw - 0.85) / 2, grid_y + 0.6, 0.85)
    text_box(s, cx + 0.22, grid_y + 1.75, cw - 0.44, 2.6,
             text="AI по такому же промпту сгенерирует похожую панораму долины "
                  "за секунды — с дальними горами и характерным тёплым светом — но "
                  "без истории создания и без оригинальной руки художника",
             size=14, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.35)

    cx2 = cx + cw + gap
    ocean_box(s, cx2, grid_y, cw, ch)
    pad2 = 0.22
    icon(s, "book-open", "028090", 64, cx2 + pad2, grid_y + pad2, 0.32)
    # NOTE: multipara_box (not literal \n in a single run) per notes/mcp-limitations.md
    # [#sem01-render-1] — \n inside one run does not reliably line-break under LibreOffice.
    multipara_box(s, cx2 + pad2 + 0.42, grid_y + pad2 - 0.02, cw - 2 * pad2 - 0.42, 0.6, [
        {"text": "Хосе Мария Веласко,", "size": 11, "bold": True, "color": MID, "line_spacing": 1.15},
        {"text": "«Долина Мехико с горы Санта-Исабель», 1875", "size": 11, "bold": True, "color": MID, "line_spacing": 1.15},
    ])
    img_h2 = ch - 1.15 - 0.32
    img_path = SHOTS / "s13-velasco-real.jpg"
    add_image(s, img_path, cx2 + pad2, grid_y + pad2 + 0.62, w=cw - 2 * pad2, h=img_h2)
    # Attribution caption BELOW the image (on card surface), not overlaid on dark
    # painting area — WCAG contrast fix (P-list finding, iter4 review, round 1).
    text_box(s, cx2 + pad2, grid_y + pad2 + 0.62 + img_h2 + 0.06, cw - 2 * pad2, 0.28,
             text="общественное достояние · Wikimedia Commons", size=9, italic=True, color=LIGHT)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Кто автор — визуализация данных: AI или человек?", size=24)

    grid_y = 1.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.9
    # Round-2 A/B shuffle: A=AI-generated chart, B=real OWID chart (flipped).
    # Round-3 point fix (student-simulator P1): the raw OWID PNG (used as-is
    # here before) still had the "Our World in Data" logo, the English title
    # "Share of the population using the Internet", and the source caption
    # baked into the image — all identifying marks that gave away which
    # chart was real before the vote. The question card (s15) now uses a
    # pre-cropped version with the title/logo band and the source-caption
    # band removed (plot area only, no attribution). The reveal slide (s16)
    # still uses the original full image with attribution — that is where
    # disclosure belongs.
    charts = ["s16-ai-generated-chart.png", "s15-owid-internet-noattrib.png"]
    for i, chart in enumerate(charts):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.22
        text_box(s, cx + pad, grid_y + pad - 0.02, cw - 2 * pad, 0.35,
                 text=f"График {chr(65+i)}", size=14, bold=True, color=MID)
        # Round-2: removed AI/человек pill chips (hand-raise voting, not
        # click) — chart image now uses the freed vertical space.
        img_path = ASSETS / "charts" / chart
        add_image(s, img_path, cx + pad, grid_y + pad + 0.4, w=cw - 2 * pad, h=ch - 0.75)
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "График A — AI, секунды. График B — Our World in Data, недели верификации", size=17)
    grid_y = 1.65
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 3.35
    # Round-2 A/B shuffle: A=AI, B=OWID (flipped).
    labels = ["Сгенерировано через AI-инструмент на вымышленных данных",
              "Our World in Data · данные ITU / World Bank · CC-BY"]
    charts = ["s16-ai-generated-chart.png", "s16-owid-internet.png"]
    for i, (chart, lbl) in enumerate(zip(charts, labels)):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        pad = 0.2
        icon(s, "sparkles" if i == 0 else "bar-chart-2", "028090", 64,
             cx + pad, grid_y + pad, 0.3)
        text_box(s, cx + pad + 0.4, grid_y + pad - 0.03, cw - 2 * pad - 0.4, 0.5,
                 text=lbl, size=10.5, bold=True, color=MID, line_spacing=1.15,
                 anchor=MSO_ANCHOR.MIDDLE)
        img_path = ASSETS / "charts" / chart if "ai-generated" in chart else SHOTS / chart
        add_image(s, img_path, cx + pad, grid_y + pad + 0.55, w=cw - 2 * pad, h=ch - 0.85)

    # Tradeoff block — round-2 rewrite: no put-down framing ("без верификации,
    # без гарантии" was read as diminishing). New framing: no authorship/ethics
    # question here (unlike painting) — the only question is data provenance.
    tw_y = grid_y + ch + 0.22
    tw_h = 7.2 - tw_y
    ocean_box(s, 0.55, tw_y, 12.23, tw_h, fill=SURFACE, stroke=TEAL)
    half = (12.23 - 0.6) / 2
    text_box(s, 0.85, tw_y + 0.14, half, 0.32, text="AI-график", size=13,
             bold=True, color=GOLD)
    text_box(s, 0.85, tw_y + 0.5, half, tw_h - 0.65,
             text="Решает задачу визуализации за секунды · здесь нет вопроса авторства или этики, как с картиной",
             size=11.5, color=DEEP, line_spacing=1.2)
    text_box(s, 0.85 + half, tw_y + 0.14, half, 0.32, text="Our World in Data", size=13,
             bold=True, color=MID)
    text_box(s, 0.85 + half, tw_y + 0.5, half, tw_h - 0.65,
             text="Исследовательская команда собирает и верифицирует источники неделями",
             size=11.5, color=DEEP, line_spacing=1.2)
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Верно или неверно? Проверим интуицию про AI", size=25)

    list_x, list_y = 0.55, 1.6
    list_w = 12.23
    n = 6
    gap = 0.14
    row_h = (5.5 - gap * (n - 1)) / n

    # Round-2: removed 2 spoiler inline glosses ("(токены, не буквы)" and
    # "(контекстное окно)") that gave away the answer before the vote; deleted
    # the "код без ошибок = логически верен" statement entirely (redundant with
    # the s18 memo's own item 3, which is also being retired); replaced the
    # context-window statement with a new determinism/training-data statement.
    rows = [
        ("AI считает буквы в слове напрямую, как человек", "hash", None),
        ("AI всегда знает, что произошло в мире сегодня", "calendar", None),
        ("AI иногда уверенно говорит неправду, не намереваясь обмануть", "brain-circuit", None),
        ("AI-чат по умолчанию видит вашу личную почту и файлы", "lock", None),
        ("AI всегда выдаёт ответ в соответствии с теми данными, которые в него заложили", "database", None),
        ("Все AI-модели дают одинаковый ответ на один и тот же вопрос", "shuffle", None),
    ]
    for i, (stmt, ic, gloss) in enumerate(rows):
        ry = list_y + i * (row_h + gap)
        ocean_box(s, list_x, ry, list_w, row_h)
        text_box(s, list_x + 0.16, ry, 0.4, row_h,
                 text=str(i + 1), size=15, bold=True, color=LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
        icon_x = list_x + 0.6
        icon(s, ic, "028090", 64, icon_x, ry + (row_h - 0.34) / 2, 0.34)
        text_x = icon_x + 0.48
        text_w = list_w - (text_x - list_x) - 2.85
        text_box(s, text_x, ry, text_w, row_h,
                 text=stmt, size=13, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        # Round-3 point fix (student-simulator, optional): flat text labels
        # instead of rounded-pill chips — pills visually echoed the
        # AI/человек pill-buttons removed from s13/s15/s16 in round-2, which
        # could read as "the buttons are back". Flat labels + a thin
        # vertical divider are a distinct shape language (no rounded outline).
        vote_x = list_x + list_w - 2.75
        text_box(s, vote_x, ry + (row_h - 0.36) / 2, 1.28, 0.36, text="верно",
                 size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        divider = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(vote_x + 1.32), Inches(ry + (row_h - 0.28) / 2),
            Inches(0.014), Inches(0.28))
        divider.fill.solid(); divider.fill.fore_color.rgb = LIGHT
        divider.line.fill.background()
        disable_shadow(divider)
        text_box(s, vote_x + 1.42, ry + (row_h - 0.36) / 2, 1.35, 0.36, text="неверно",
                 size=12, bold=True, color=SLATE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, list_x, list_y + 5.5 + 0.08, list_w, 0.28,
             text="голосуем до объяснения", size=11, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s17"))


def build_s18(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Памятка «AI reality-check» — 5 пунктов на неделю вперёд", size=25)

    # Round-2: item 3 ("код без ошибок...") removed — referred to the deleted
    # quiz statement. New item added on determinism/training-data (mirrors the
    # new s17 statement 5). All items now render with identical fill/stroke —
    # the round-1 gold highlight + "самая частая ошибка" tag on item 1 removed
    # per brief (all 5 items should look equal).
    # Round-3: item 4 referred to the context-window quiz statement that was
    # already deleted in round-2 (orphan reference) — replaced with a new
    # item on tokens, mirroring quiz statement 1 (previously uncovered by
    # the memo).
    items = [
        "Уверенный тон ответа AI не означает, что ответ фактически верен",
        "AI не знает о событиях после даты среза обучения без веб-поиска",
        "Ответ AI не выводится однозначно из тренировочных данных — из-за "
        "случайности генерации один и тот же вопрос может дать разные ответы",
        "AI видит текст не буквами, а токенами — поэтому иногда путается "
        "в побуквенных задачах вроде подсчёта букв в слове",
        "Разные AI-модели могут давать разные ответы — сверка с несколькими источниками разумна",
    ]
    list_x, list_y = 0.55, 1.7
    list_w = 12.23
    n = len(items)
    gap = 0.18
    row_h = (5.3 - gap * (n - 1)) / n
    for i, txt in enumerate(items):
        ry = list_y + i * (row_h + gap)
        ocean_box(s, list_x, ry, list_w, row_h, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5)
        badge = filled_rect(s, list_x + 0.2, ry + (row_h - 0.5) / 2, 0.5, 0.5,
                             LIGHT, radius=True, radius_adj=0.5)
        text_box(s, list_x + 0.2, ry + (row_h - 0.5) / 2, 0.5, 0.5,
                 text=str(i + 1), size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, list_x + 0.95, ry, list_w - 1.15, row_h,
                 text=txt, size=14.5, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    hero_w = 6.283
    # highlight_idx=4 -> "Разработка / IT" (datacenter tile), closest theme to Lecture 1
    industry_photo_grid(s, 0.0, 0.0, hero_w, 7.5, highlight_idx=4)

    rx = hero_w + 0.55
    rw = 13.333 - rx - 0.55
    text_box(s, rx, 1.35, rw, 0.5,
             text="СПАСИБО", size=14, bold=True, color=TEAL)
    text_box(s, rx, 1.85, rw, 2.3,
             text="От сегодняшнего среза — к системному разбору AI по отраслям",
             size=27, bold=True, color=DEEP, line_spacing=1.15)
    text_box(s, rx, 4.1, rw, 1.4,
             text="Сегодня — ваш личный опыт. Дальше — Лекция 1 и системный разбор, "
                  "отрасль за отраслью.",
             size=15.5, color=MID, line_spacing=1.35)
    chip(s, rx, 5.6, 3.3, 0.55, "Лекция 1  →  далее", fill=DEEP, size=14)
    speaker_notes(s, load_notes("s19"))


# ============================================================
# Main
# ============================================================
# NOTE: this list still reflects the OLD 19-slide / pre-bio-slide numbering.
# See the "KNOWN DRIFT" block in the module docstring above — running this
# script as-is regenerates the OLD deck, not the current 20-slide
# rendered/sem-01.pptx. Do not run until the TODO items there are done.
BUILDERS = [
    build_s01, build_s02, build_s03, build_s04, build_s05,
    build_s06, build_s07, build_s08, build_s09, build_s10,
    build_s11, build_s12, build_s13, build_s14, build_s15,
    build_s16, build_s17, build_s18, build_s19,
]


def main():
    p = setup_pres()
    for fn in BUILDERS:
        fn(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved {OUT} ({len(BUILDERS)} slides)")


if __name__ == "__main__":
    main()
