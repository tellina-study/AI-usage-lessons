"""
apply_comments_r5.py — applies editor comments #249-#256 to lec-01.pptx IN PLACE.

build_lec01.py is NOT run (its main() would rebuild). We import its styling helpers
and palette only, and edit the CURRENT rendered/lec-01.pptx directly.

Comments:
  #249 slide 10 — add 1956 Dartmouth ("рождение термина AI") to the "Открытия" row
  #250 slide 12 — refresh ChatGPT WAU stat 900M (фев 2026) -> 1B (авг 2026)
  #251 swap slides 15 <-> 16
  #252 delete slide 20 (ЧАТ — кейс)
  #253 slide 22 — visually mark inner (1 file) + outer (x200 files) loops
  #254 slide 30 — AGI table verified current (no change; see reply)
  #255 slide 34 — grading: drop 10 attendance, exam 30 -> 40
  #256 move semester-map + grading to front, insert MAX chat page after grading,
       add MAX QR to final Q&A slide
"""
import copy
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import build_lec01 as B
from build_lec01 import (
    text_box, ocean_box, filled_rect, chip, add_image, set_slide_bg,
    speaker_notes, disable_shadow,
    DEEP, MID, LIGHT, TEAL, SURFACE, WHITE, GOLD, SLATE, GOLD_TINT, SOFT_GREY,
)

HERE = Path(__file__).resolve().parent
PPTX = HERE / "lec-01.pptx"
SEM01 = HERE / "../../../seminars/sem-01/rendered/sem-01.pptx"
QR = HERE / "assets/max-qr.png"

prs = Presentation(str(PPTX))
S = prs.slides


def emu_in(v):
    return Emu(v).inches if v is not None else 0.0


def set_lines(shp, lines):
    paras = shp.text_frame.paragraphs
    for i, line in enumerate(lines):
        if i < len(paras) and paras[i].runs:
            paras[i].runs[0].text = line
            for r in paras[i].runs[1:]:
                r.text = ''


def set_only_run_text(shp, text):
    shp.text_frame.paragraphs[0].runs[0].text = text


def del_shape(shp):
    shp._element.getparent().remove(shp._element)


def shape_by_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


# ============================================================
# #249 — slide 10: reflow "Открытия" row 3 -> 4 items (add 1956 Dartmouth)
# ============================================================
def fix_s10():
    s = S[9]
    def find(txt):
        return shape_by_text(s, txt)
    t_turing = find('Тьюринг')
    y_1950 = find('1950')
    t_eliza = find('ELIZA')
    y_1966 = find('1966')
    t_expert = find('Экспертные системы')
    y_1980 = find('1980-е')
    # ticks in the "Открытия" band (top ~2.69): 3 small rectangles
    ticks = [sh for sh in s.shapes if sh.shape_type == 1 and abs(emu_in(sh.top) - 2.69) < 0.05
             and emu_in(sh.width) < 0.2]
    ticks.sort(key=lambda sh: emu_in(sh.left))

    W = 2.72
    xs = [1.00, 3.85, 6.70, 9.55]
    centers = [x + W / 2 for x in xs]

    # place existing 3 into columns 0,2,3 ; new Dartmouth into col 1
    title_map = [(t_turing, 0), (t_eliza, 2), (t_expert, 3)]
    year_map = [(y_1950, 0), (y_1966, 2), (y_1980, 3)]
    for sh, col in title_map:
        sh.left = Inches(xs[col]); sh.width = Inches(W); sh.top = Inches(2.05); sh.height = Inches(0.72)
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)
        sh.text_frame.word_wrap = True
    for sh, col in year_map:
        sh.left = Inches(xs[col]); sh.width = Inches(W)
    # reposition existing ticks to centers 0,2,3
    for sh, col in zip(ticks, [0, 2, 3]):
        sh.left = Inches(centers[col] - emu_in(sh.width) / 2)

    # new column 1 — deep-copy title, year, tick from Turing's
    new_title_el = copy.deepcopy(t_turing._element)
    s.shapes._spTree.append(new_title_el)
    nt = s.shapes[-1]
    set_lines(nt, ['Дартмутский семинар — рождение термина «AI»'])
    nt.left = Inches(xs[1]); nt.width = Inches(W); nt.top = Inches(2.05); nt.height = Inches(0.72)
    for p in nt.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11)
    nt.text_frame.word_wrap = True

    new_year_el = copy.deepcopy(y_1950._element)
    s.shapes._spTree.append(new_year_el)
    ny = s.shapes[-1]
    set_lines(ny, ['1956'])
    ny.left = Inches(xs[1]); ny.width = Inches(W)

    if ticks:
        new_tick_el = copy.deepcopy(ticks[0]._element)
        s.shapes._spTree.append(new_tick_el)
        ntk = s.shapes[-1]
        ntk.left = Inches(centers[1] - emu_in(ntk.width) / 2)
    print('  #249 s10: 4-column Открытия row (added 1956 Дартмут)')


# ============================================================
# #250 — slide 12: refresh ChatGPT WAU stat
# ============================================================
def fix_s12():
    s = S[11]
    changed = []
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == '900M':
            set_only_run_text(sh, '1B'); changed.append('900M->1B')
        elif t == 'ChatGPT, февраль 2026':
            set_only_run_text(sh, 'ChatGPT, август 2026'); changed.append('date')
    # update the headline number reference too
    head = shape_by_text(s, '900M пользователей')
    if head:
        for p in head.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('900M пользователей', '~1 млрд пользователей')
        changed.append('headline')
    print('  #250 s12:', changed)


# ============================================================
# #253 — slide 22: visualize inner (1 file) + outer (x200) loops
# ============================================================
def _dash(line):
    ln = line._get_or_add_ln()
    d = ln.find(qn('a:prstDash'))
    if d is None:
        d = ln.makeelement(qn('a:prstDash'), {})
        ln.append(d)
    d.set('val', 'dash')


def fix_s22():
    s = S[21]
    # Inner-loop box around steps 2..6 (ovals at top 3.17 .. 5.53)
    x, y, w, h = 5.03, 3.05, 7.66, 2.85
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.color.rgb = TEAL
    box.line.width = Pt(1.6)
    _dash(box.line)
    disable_shadow(box)
    try:
        box.adjustments[0] = 0.05
    except Exception:
        pass
    # inner-loop label — vertical, in the right margin alongside the box
    lab = s.shapes.add_textbox(Inches(11.75), Inches(3.95), Inches(2.6), Inches(0.4))
    lab.rotation = 270
    ltf = lab.text_frame
    ltf.word_wrap = False
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = 'цикл: обработка одного файла'
    lr.font.size = Pt(11)
    lr.font.bold = True
    lr.font.color.rgb = TEAL
    # outer-loop return arrow: elbow in the narrow left gap from step 7 up to step 2
    ax = 4.965  # panel left edge ~4.90, ovals at 5.20
    bar = filled_rect(s, ax, 3.30, 0.055, 3.00, TEAL)  # vertical spine
    disable_shadow(bar)
    stub = filled_rect(s, ax, 6.25, 0.30, 0.055, TEAL)  # bottom stub toward step7
    disable_shadow(stub)
    # up-arrow head into step 2
    head = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ax - 0.06), Inches(3.10), Inches(0.18), Inches(0.22))
    head.rotation = 0
    head.fill.solid(); head.fill.fore_color.rgb = TEAL
    head.line.fill.background()
    disable_shadow(head)
    # outer-loop label chip below step 7
    chip(s, 5.20, 6.46, 2.55, 0.32, '↻  повторить × 200 файлов',
         fill=GOLD, color=DEEP, size=11, bold=True)
    print('  #253 s22: inner dashed box + outer return arrow + labels')


# ============================================================
# #255 — slide 34: grading formula
# ============================================================
def fix_s34():
    s = S[33]
    sh = shape_by_text(s, 'посещаемость')
    runs = sh.text_frame.paragraphs[0].runs
    # runs: 0='100'(gold) 1='  =  10 ' 2='(посещаемость)' 3='  +  30 ' 4='(экзамен)' 5='  +  3×20 ' 6='(РК..)'
    runs[1].text = '  =  40 '
    runs[2].text = '(экзамен)'
    runs[3].text = ''
    runs[4].text = ''
    print('  #255 s34: grading formula -> 40 exam, no attendance')


# ============================================================
# #256 support — cross-file copy of the sem-01 MAX chat slide
# ============================================================
def copy_max_slide():
    src_prs = Presentation(str(SEM01))
    src = src_prs.slides[4]  # s04b MAX chat
    dest = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    # background
    src_csld = src._element.find(qn('p:cSld'))
    dst_csld = dest._element.find(qn('p:cSld'))
    bg = src_csld.find(qn('p:bg'))
    if bg is not None:
        dst_csld.insert(0, copy.deepcopy(bg))
    else:
        set_slide_bg(dest, WHITE)
    # copy shapes
    for shp in src.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    # remap image rels (cross-package: create new image parts in dest)
    for pic in dest.shapes._spTree.iter(qn('p:pic')):
        blip = pic.find('.//' + qn('a:blip'))
        if blip is None:
            continue
        rId = blip.get(qn('r:embed'))
        if not rId:
            continue
        src_img = src.part.related_part(rId)
        img_part, new_rId = dest.part.get_or_add_image_part(BytesIO(src_img.blob))
        blip.set(qn('r:embed'), new_rId)
    # fresh notes
    speaker_notes(dest, MAX_NOTE)
    return list(prs.slides._sldIdLst)[-1]


MAX_NOTE = (
    "Всё общение по курсу — в чате в мессенджере MAX. Отсканируйте QR-код камерой телефона "
    "или перейдите по ссылке max.ru/join/sHoHlhI4jvW_swcq5ZLowvz6G94m0IMhAwngxuxsEpI "
    "и вступите в чат до первой лекции.\n\n"
    "Что там будет. Во-первых, объявления: расписание, любые изменения, важные новости — "
    "всё публикуется в чате, отдельно письма ждать не нужно. Во-вторых, материалы: слайды "
    "лекций и семинаров я выкладываю туда по мере проведения занятий, так что всё собрано в "
    "одном месте. В-третьих, вопросы: если что-то непонятно между парами — спрашивайте в "
    "чате, не обязательно копить до семинара.\n\n"
    "Пожалуйста, вступите сегодня — так вы точно не пропустите первое объявление."
)


# ============================================================
# #256 — add MAX QR to final Q&A slide
# ============================================================
def add_qr_to_qa():
    s = S[35]  # Q&A (Вопросы?)
    # small card bottom-right
    cx, cy, cw, ch = 9.7, 4.55, 3.1, 2.6
    ocean_box(s, cx, cy, cw, ch, fill=WHITE, stroke=LIGHT, stroke_pt=1.6)
    qsz = 1.55
    add_image(s, str(QR), cx + (cw - qsz) / 2, cy + 0.22, w=qsz, h=qsz)
    text_box(s, cx + 0.15, cy + 1.86, cw - 0.3, 0.32, 'Чат курса в MAX',
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, cx + 0.12, cy + 2.16, cw - 0.24, 0.4,
             'max.ru/join/sHoHlhI4jvW…', size=9.5, color=MID, align=PP_ALIGN.CENTER)
    print('  #256 s36: MAX QR added to Q&A slide')


# ============================================================
# Orchestrate
# ============================================================
print('Content edits:')
fix_s10()
fix_s12()
fix_s22()
fix_s34()
add_qr_to_qa()

print('Structural edits:')
lst = prs.slides._sldIdLst
ids = list(lst)  # 36 original sldId elements (order = slides 1..36)
assert len(ids) == 36, f'expected 36, got {len(ids)}'

chat_id = copy_max_slide()  # appended -> now 37 in list
print('  #256: MAX chat slide copied from sem-01')

# indices (0-based): 0..5 intro(s01..s06); 6..31 middle(s07..s32);
# 14=s15,15=s16 (swap #251); 19=s20 (delete #252); 32=s33 map,33=s34 grading (move #256);
# 34=s35 teaser, 35=s36 Q&A
head_intro = ids[0:6]
moved = [ids[32], ids[33]]          # map, grading -> front
tail = [ids[34], ids[35]]           # teaser, Q&A

middle = ids[6:32]                   # s07..s32 (26 elements)
# swap s15<->s16 : positions 8 and 9 within middle (global 14,15)
middle[8], middle[9] = middle[9], middle[8]
# delete s20 : global 19 -> position 13 within middle
del_el = ids[19]
middle = [e for e in middle if e is not del_el]

desired = head_intro + moved + [chat_id] + middle + tail
assert len(desired) == 36, f'desired {len(desired)}'

# drop the deleted slide's relationship + part
del_rId = del_el.get(qn('r:id'))
prs.part.drop_rel(del_rId)
print('  #252: slide 20 (ЧАТ кейс) deleted')
print('  #251: slides 15<->16 swapped')
print('  #256: map+grading moved to front, chat inserted after grading')

for el in list(lst):
    lst.remove(el)
for el in desired:
    lst.append(el)

prs.save(str(PPTX))
print(f'Saved {PPTX} — total slides now {len(list(prs.slides))}')
