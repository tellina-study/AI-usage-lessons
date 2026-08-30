"""
Full 33-slide build of Лекции 1 v3.1 (Phase 12.4 revision of EPIC #64, issue #70).

Source-of-truth: deck.yaml v3.1 + chapter v3.1 (status=reviewed, 16406 слов) +
slides/*.md v3.1 (33 файлов с readable speaker notes 150-300 слов).

v3.1 changes vs v3 (4-critic synthesis 2026-05-13):
- 34 → 33 slides: removed s26 ARC-AGI economics (concept retains poorly при self-study)
  and s28 Pearl 3 levels (концептуально красиво, но к концу 75-min лекции не заходит).
- Renumbered: s27→s26 (4-speaker AGI table), s29→s28 (summary+homework), s30→s29 (roadmap),
  s31→s30 (lec2 teaser), s32→s31 (Q&A).
- Added NEW s27 — section 5 divider («Что забрать домой») per DoD §10 + reader-rendered feedback.
- Critical fixes (4):
  * s13 speaker notes synced with visual (Model = left-top, Agent = right-bottom).
  * «Приложение-робот» renamed to «Приложение (автоматизация)» on s21 quadrant;
    s20+s21 notes explain 2 types of apps (with UI / without UI).
  * s05b funnel «10% в проде» → «10% доходят до прода»; widened gold plate.
  * NEW s27 divider section 5.
- High-value fixes (7):
  * s13 axis labels enlarged (10pt → 13-16pt).
  * s15 RU/EN sub-labels unified to RU.
  * s21 axis labels Q1/Q2 moved INSIDE quadrant.
  * s08 «90% откатов» n=50 caveat added to speaker notes.
  * s07 Vaswani citation timestamp «на май 2026».
  * s29 PARTS disclaimer added to speaker notes.
  * s28 takeaway 3 — removed Pearl reference (Pearl slide deleted).

v3 baseline (preserved): Ocean Gradient palette, Ocean rounded box motif on every content
slide, Gold ≥1×/slide, footer-tax = 0, all notes are readable text 150-300 words.

v3.3 (issue #153, 21-fix polish): 34 slides — see deck.yaml v3_3_changes for detail.

v3.4 (issue #155 batch 1, 10-fix polish + section-divider audit): 36 slides —
added s05c (section 1 divider) + s07a (section 2 divider), both reusing the
existing nav_slide() unified template (same as s10/s22/s27). See deck.yaml
v3_4_changes for full fix-by-fix mapping. chapter.md/speech.md NOT yet synced
(deferred until all 4 issue #155 batches land).

Canvas: 13.333" × 7.5" (16:9). Pacing (v3.4): 61.5 active + 13.5 buffer = 75 min.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parents[1]      # library/lectures/lec-01 (worktree-relative)
ASSETS = ROOT / "rendered/assets"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-01.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


# ============================================================
# Reference system (ported from lec-04 _helpers.py, issue #170).
# Small superscript muted [N] markers inside body text + a bottom clickable
# numbered source list + «Источники:» block in speaker notes. All URLs come
# ONLY from URLS (reference-registry.md); volatile → [VFY-day-of] in notes.
# ============================================================
_REF_RE = re.compile(r'\[\d+(?:\s*[,–—-]\s*\d+)*\]')
_AMAIN = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _run_props(src_run):
    f = src_run.font
    sz = f.size
    return {
        "name": f.name,
        "size_pt": (sz.pt if sz is not None else None),
        "bold": f.bold,
        "italic": f.italic,
        "color": (f.color.rgb if (f.color and f.color.type is not None) else None),
    }


def _clone_run_after(anchor_r, props, text, *, ref=False,
                     ref_frac=0.52, ref_color=None):
    """Insert a new <a:r> right after anchor_r with cloned props (or a small
    superscript muted variant when ref=True)."""
    if ref_color is None:
        ref_color = LIGHT
    new_r = etree.SubElement(anchor_r.getparent(), _AMAIN + "r")
    anchor_r.addnext(new_r)
    rpr = etree.SubElement(new_r, _AMAIN + "rPr")
    base = props["size_pt"] or 16.0
    if ref:
        rpr.set("sz", str(int(round(base * ref_frac * 100))))
        rpr.set("baseline", "30000")
        rpr.set("b", "0")
        rpr.set("i", "1")
    else:
        if props["size_pt"] is not None:
            rpr.set("sz", str(int(round(base * 100))))
        if props["bold"] is not None:
            rpr.set("b", "1" if props["bold"] else "0")
        if props["italic"] is not None:
            rpr.set("i", "1" if props["italic"] else "0")
    if props["name"]:
        for tag in ("latin", "cs", "ea"):
            el = etree.SubElement(rpr, _AMAIN + tag)
            el.set("typeface", props["name"])
    col = ref_color if ref else props["color"]
    if col is not None:
        fill = etree.SubElement(rpr, _AMAIN + "solidFill")
        clr = etree.SubElement(fill, _AMAIN + "srgbClr")
        clr.set("val", str(col))
    t = etree.SubElement(new_r, _AMAIN + "t")
    t.text = text
    return new_r


def shrink_refs_in_frame(text_frame, *, ref_frac=0.52, ref_color=None):
    """Split every [N] marker inside the frame into a small superscript muted
    run. Non-destructive to surrounding text formatting."""
    if ref_color is None:
        ref_color = LIGHT
    for para in text_frame.paragraphs:
        for run in list(para.runs):
            txt = run.text
            if not txt or "[" not in txt:
                continue
            matches = list(_REF_RE.finditer(txt))
            if not matches:
                continue
            props = _run_props(run)
            run.text = txt[:matches[0].start()]
            anchor = run._r
            for i, m in enumerate(matches):
                anchor = _clone_run_after(anchor, props, m.group(),
                                          ref=True, ref_frac=ref_frac,
                                          ref_color=ref_color)
                nxt = matches[i + 1].start() if i + 1 < len(matches) else len(txt)
                between = txt[m.end():nxt]
                if between:
                    anchor = _clone_run_after(anchor, props, between, ref=False)
    return text_frame


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    shrink_refs_in_frame(tf)
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    shrink_refs_in_frame(tf)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0, radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    elif h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.45, h=1.15, w=12.3, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.18, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True):
    box = filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                      radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.08, w=w - 0.4, h=h - 0.16, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


def speaker_notes(slide, text):
    """Write notes as readable PARAGRAPHS (lec-04 pattern): split on blank
    lines → one notes-paragraph each; the «Источники:» block keeps its own
    hard line breaks so each [N] URL sits on its own line."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    blocks = [b.strip() for b in re.split(r'\n\s*\n', text.strip()) if b.strip()]
    if not blocks:
        blocks = [""]
    for i, block in enumerate(blocks):
        if block.lstrip().startswith("Источники:"):
            lines = [ln.rstrip() for ln in block.split("\n")]
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = lines[0]
            for ln in lines[1:]:
                sub = tf.add_paragraph()
                sub.text = ln
            continue
        one = re.sub(r'\s*\n\s*', ' ', block)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = one


def eyebrow_pill(slide, text):
    """Consistent eyebrow pill top-left — issue #153 fix #10.

    Marks which of the 4 implementation types (МОДЕЛЬ / ЧАТ / АГЕНТ) a
    content slide discusses. Used on s15/s16/s17/s18/s19/s19a so the
    reader always knows which layer is on screen without re-reading the
    whole assertion.
    """
    w = 0.35 + 0.16 * len(text)
    h = 0.42
    x, y = 0.55, 0.35
    chip(slide, x, y, w, h, text, fill=DEEP, color=WHITE, size=13, bold=True)


# ============================================================
# Speaker notes loader from md
# ============================================================
def load_notes(slide_id):
    """Extract Speaker notes block from slide markdown.

    Fix-1 (2026-05-13): Лектору block removed from all slide MD — speaker
    notes now contain ONLY the readable student-facing text (150-300 words).
    Lecturer-side cues live in speech.md, not in slide notes.
    """
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    notes_match = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)', md, re.DOTALL)
    notes = notes_match.group(1).strip() if notes_match else ""
    # Strip any trailing horizontal rule
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Canonical URL registry (issue #170) — keyed by short id.
# Source: notes/research/lecture-1/reference-registry.md §URLS.
# Only URLs from that map are used. Volatile ones get [VFY-day-of] in notes.
# Флаг 1: github_octoverse URL is canonical, but the «46% кода» claim is NOT
# confirmed by that report → the claim carries [VFY] in the s08 notes and the
# ref is NOT presented as proof of the 46% figure specifically.
# ============================================================
URLS = {
    # --- s00b / s08 macro-context ---
    "gartner_2024": "https://www.gartner.com/en/newsroom/press-releases/2024-10-03-gartner-says-generative-ai-will-require-80-percent-of-engineering-workforce-to-upskill-through-2027",
    "cnews_vedomosti": "https://www.vedomosti.ru/technology/articles/2026/03/24/1184974-biznes-svernul-ili-otlozhil-9-iz-10-proektov-po-vnedreniyu-generativnogo-ii",
    # --- s01 demo ---
    "yolov8": "https://docs.ultralytics.com/models/yolov8/",
    "mediapipe": "https://ai.google.dev/edge/mediapipe/solutions/guide",
    # --- s06 / s07 / s12 definitions & history ---
    "aima": "https://aima.cs.berkeley.edu/",
    "iso_22989": "https://www.iso.org/standard/74296.html",
    "mitchell": "https://www.cs.cmu.edu/~tom/mlbook.html",
    "mccorduck": "https://www.google.com/books/edition/Machines_Who_Think/dPGij4vsHKgC",
    "mcculloch_pitts": "https://doi.org/10.1007/BF02478259",
    "vaswani": "https://arxiv.org/abs/1706.03762",
    "dhar": "https://doi.org/10.1145/3664804",
    "goodfellow": "https://www.deeplearningbook.org/",
    # --- s08 scale ---
    "so_2025": "https://survey.stackoverflow.co/2025",
    "openai_wau": "https://techcrunch.com/2026/02/27/chatgpt-reaches-900m-weekly-active-users/",
    "github_octoverse": "https://github.blog/news-insights/octoverse/octoverse-a-new-developer-joins-github-every-second-as-ai-leads-typescript-to-1/",
    "gvr": "https://www.grandviewresearch.com/industry-analysis/artificial-intelligence-ai-market",
    # --- s09 breakthroughs ---
    "mistral_7b": "https://mistral.ai/news/announcing-mistral-7b",
    "deepseek_r1": "https://arxiv.org/abs/2501.12948",
    "semianalysis": "https://semianalysis.com/2025/01/31/deepseek-debates/",
    "bloomberg_deepseek": "https://www.bloomberg.com/news/articles/2025-01-27/asml-sinks-as-china-ai-startup-triggers-panic-in-tech-stocks",
    "openclaw": "https://github.com/openclaw/openclaw",
    "llamacpp": "https://github.com/ggml-org/llama.cpp",
    # --- s11 / s13 / s16 / s18 / s20 / s21 agents ---
    "anthropic_agents": "https://www.anthropic.com/research/building-effective-agents",
    "weng": "https://lilianweng.github.io/posts/2023-06-23-agent/",
    "mcp": "https://www.anthropic.com/news/model-context-protocol",
    "react": "https://arxiv.org/abs/2210.03629",
    "autonomy_levels": "https://arxiv.org/abs/2506.12469",
    "google_agents_wp": "https://www.kaggle.com/whitepaper-agents",
    "ng_patterns": "https://www.deeplearning.ai/the-batch/how-agents-can-improve-llm-performance/",
    # --- s15 model ---
    "kreuzberger": "https://arxiv.org/abs/2205.02302",
    "alphafold": "https://doi.org/10.1038/s41586-021-03819-2",
    # --- s17 chat ---
    "vciom": "https://wciom.ru/analytical-reviews/analiticheskii-obzor/neiroseti-v-nashei-zhizni",
    "dam": "https://arxiv.org/abs/2406.16937",
    # --- s20 application ---
    "google_translate": "https://blog.google/products-and-platforms/products/translate/fun-facts-google-translate-20-years/",
    # --- s22 / s23 boundaries & governance ---
    "nist_rmf": "https://www.nist.gov/itl/ai-risk-management-framework",
    "nist_600": "https://nvlpubs.nist.gov/nistpubs/ai/NIST.AI.600-1.pdf",
    "eu_ai_act": "https://eur-lex.europa.eu/eli/reg/2024/1689/oj/eng",
    "bloomberg_samsung": "https://www.bloomberg.com/news/articles/2023-05-02/samsung-bans-chatgpt-and-other-generative-ai-use-by-staff-after-leak",
    "openai_terms": "https://openai.com/enterprise-privacy/",
    # --- s24 hallucinations ---
    "huang": "https://arxiv.org/abs/2311.05232",
    "ji": "https://doi.org/10.1145/3571730",
    "vectara": "https://github.com/vectara/hallucination-leaderboard",
    "cybsafe": "https://www.staysafeonline.org/articles/oh-behave-the-annual-cybersecurity-attitudes-and-behaviors-report-2025",
    # --- s25 bias/sycophancy ---
    "sycophancy": "https://openai.com/index/sycophancy-in-gpt-4o/",
    "reward_misspec": "https://arxiv.org/abs/2201.03544",
    # --- s26 AGI ---
    "searle": "https://doi.org/10.1017/S0140525X00005756",
    "bostrom": "https://global.oup.com/academic/product/superintelligence-9780199678112",
}


# ============================================================
# Per-slide source registry — drives BOTH the bottom clickable [N] list AND
# the «Источники:» block in the speaker notes, so slide-[N] and notes-[N] can
# never diverge. Entry: (num, short_name, urlkey, gloss[, volatile]).
# volatile → «[VFY-day-of]» appended in notes only.
# ============================================================
SLIDE_REFS = {
    "s00b": [
        ("1", "Gartner — пресс-релиз (окт 2024)", "gartner_2024",
         "80% инженерных работников должны осваивать GenAI к 2027", True),
        ("2", "Vedomosti / Intellectual Analytics (мар 2026)", "cnews_vedomosti",
         "9 из 10 корпоративных GenAI-пилотов в РФ свёрнуты/отложены (n=50 крупных орг.)", True),
    ],
    "s01": [
        ("1", "Ultralytics — YOLOv8 (2023)", "yolov8",
         "narrow-модель детекции; локальный inference ~30 fps на CPU без интернета"),
        ("2", "Google — MediaPipe", "mediapipe",
         "on-device real-time ML pipelines — класс локальных narrow-решений"),
    ],
    "s06": [
        ("1", "Russell & Norvig — AIMA, 4-е изд. (2021)", "aima",
         "4 определения по 2 осям (думать/действовать × человекоподобно/рационально)"),
        ("2", "ISO/IEC 22989:2022", "iso_22989",
         "AI-система = engineered system для человеко-заданных целей; опора EU AI Act"),
        ("3", "Mitchell — Machine Learning (1997)", "mitchell",
         "функциональное определение: поведение из обученной модели = AI"),
        ("4", "Searle — Chinese Room (1980) / Turing (1950)", "searle",
         "AI Effect + возражение Сёрла: поведение на бенчмарке ≠ понимание"),
    ],
    "s06a": [
        ("1", "McCulloch & Pitts (1943)", "mcculloch_pitts",
         "формальный нейрон как логический элемент — на 13 лет раньше термина «AI»"),
    ],
    "s07": [
        ("1", "Vaswani et al. — Attention Is All You Need (2017)", "vaswani",
         "Трансформер + self-attention; >160K цитирований (май 2026) — точка перелома", True),
        ("2", "McCorduck — Machines Who Think (2004)", "mccorduck",
         "AI Effect через исторические примеры остывания задач до «просто функции»"),
        ("3", "Dhar — Paradigm Shifts in AI, CACM (2024)", "dhar",
         "смена парадигм AI как рамка для чтения таймлайна"),
    ],
    "s08": [
        ("1", "Stack Overflow — Developer Survey 2025", "so_2025",
         "n=49k+/177 стран; 51% профи ежедневно, 84% используют/планируют, 46% не доверяют коду"),
        ("2", "OpenAI — ChatGPT WAU (фев 2026)", "openai_wau",
         "~900M weekly active users — AI как инфраструктура", True),
        ("3", "GitHub — Octoverse 2025", "github_octoverse",
         "adoption Copilot; ВНИМАНИЕ: «46% кода» этим отчётом НЕ подтверждается [VFY]", True),
        ("4", "Grand View Research (2026)", "gvr",
         "AI-рынок $390.9B (2025) → $539.5B (2026), CAGR 30.6%", True),
        ("5", "Vedomosti / Intellectual Analytics (мар 2026)", "cnews_vedomosti",
         "контр-факт: ~9 из 10 пилотов РФ не доходят до прода (n=50)", True),
        ("6", "Gartner — пресс-релиз (окт 2024)", "gartner_2024",
         "80% инженеров осваивают GenAI к 2027", True),
    ],
    "s09": [
        ("1", "Mistral AI — Announcing Mistral 7B (сен 2023)", "mistral_7b",
         "Apache 2.0, обходит Llama-2 13B — малая команда уровня лидеров"),
        ("2", "DeepSeek-R1 — тех.отчёт (янв 2025)", "deepseek_r1",
         "reasoning уровня o1; 97.3% MATH-500; спорная себестоимость"),
        ("3", "SemiAnalysis — DeepSeek cost analysis (2025)", "semianalysis",
         "полная инфра $1.3–1.6B vs marginal train run V3 $5.6M — разные числа", True),
        ("4", "Bloomberg — Nvidia $589B drop (27 янв 2025)", "bloomberg_deepseek",
         "крупнейшая однодневная потеря капитализации в истории", True),
        ("5", "Steinberger — OpenClaw (GitHub)", "openclaw",
         "соло open-source агент, 100K★ за квартал; rename-churn", True),
        ("6", "Gerganov — llama.cpp / ggml.ai", "llamacpp",
         "соло→HF (фев 2026), 100K+★ быстрее PyTorch/TensorFlow", True),
    ],
    "s11": [
        ("1", "Anthropic — Building Effective Agents (2024)", "anthropic_agents",
         "слоистая эскалация: простейшее решение → наращивать при необходимости"),
        ("2", "Weng — LLM Powered Autonomous Agents (2023)", "weng",
         "Agent = LLM + Memory + Planning + Tool Use — верх слоистой модели"),
    ],
    "s12": [
        ("1", "Russell & Norvig — AIMA (2021)", "aima",
         "классификация AI-систем как рабочий язык курса"),
        ("2", "Goodfellow, Bengio, Courville — Deep Learning (2016)", "goodfellow",
         "модальности и типы задач в терминах DL"),
    ],
    "s13": [
        ("1", "Anthropic — Building Effective Agents (2024)", "anthropic_agents",
         "распределение контроля разработчик↔пользователь по мере делегирования"),
    ],
    "s15": [
        ("1", "Kreuzberger et al. — MLOps (2023)", "kreuzberger",
         "pre/post-processing вокруг модели — ответственность разработчика"),
        ("2", "Jumper et al. — AlphaFold, Nature (2021)", "alphafold",
         "канонический пример модели-прогноза; Нобель по химии 2024"),
    ],
    "s16": [
        ("1", "Anthropic — Building Effective Agents (2024)", "anthropic_agents",
         "системный промпт как инженерный рычаг; контекст собирается заново каждый шаг"),
    ],
    "s17": [
        ("1", "ВЦИОМ — «Нейросети в жизни россиян» (окт 2025)", "vciom",
         "проникновение AI-чатов в РФ; 51% пользуются раз в неделю+", True),
        ("2", "Dam et al. — Survey on LLM-based AI Chatbots (2024)", "dam",
         "таксономия LLM-чатов; чистые чаты редки в проде (расширены до агентов)"),
    ],
    "s18": [
        ("1", "Weng — LLM Powered Autonomous Agents (2023)", "weng",
         "Agent = LLM + Memory + Planning + Tool Use"),
        ("2", "Anthropic — Building Effective Agents (2024)", "anthropic_agents",
         "оркестратор + инструменты + внешняя память как слой над чатом"),
        ("3", "Anthropic — Model Context Protocol (2024)", "mcp",
         "стандарт подключения инструментов/данных к агенту"),
    ],
    "s19": [
        ("1", "Yao et al. — ReAct (2022)", "react",
         "Reasoning+Acting: план→действие→наблюдение→рефлексия с явным инструментом на шаге"),
    ],
    "s19a": [
        ("1", "Feng, McDonald, Zhang — Levels of Autonomy (2025)", "autonomy_levels",
         "5 уровней по роли пользователя: operator→collaborator→consultant→approver→observer"),
    ],
    "s20": [
        ("1", "Google — Translate at 20 (2026)", "google_translate",
         "1B+ польз./мес, ~1T слов/мес across Translate/Search/Lens — AI как функция", True),
        ("2", "Anthropic — Building Effective Agents (2024)", "anthropic_agents",
         "приложение как внешний слой: промпты скрыты, детерминированный UI"),
    ],
    "s21": [
        ("1", "Anthropic — Building Effective Agents (2024)", "anthropic_agents",
         "выбор типа реализации по 2 осям: взаимодействие × инструменты"),
        ("2", "Google — AI Agents Whitepaper (2024)", "google_agents_wp",
         "Model + Tools + Orchestration Layer — рамка квадранта"),
        ("3", "Ng — Four Agentic Design Patterns (2024)", "ng_patterns",
         "паттерны Reflection/Tool Use/Planning/Multi-Agent — когда нужен агент"),
    ],
    "s22": [
        ("1", "NIST — AI RMF 1.0 (2023)", "nist_rmf",
         "рамка управления рисками AI — граница ответственности инженера"),
        ("2", "NIST — Generative AI Profile 600-1 (2024)", "nist_600",
         "профиль рисков GenAI"),
        ("3", "EU AI Act — Reg. (EU) 2024/1689", "eu_ai_act",
         "регуляторная рамка; штрафы до 35M€/7%"),
    ],
    "s23": [
        ("1", "Bloomberg — Samsung bans ChatGPT (май 2023)", "bloomberg_samsung",
         "3 утечки за месяц → запрет внешнего GenAI; данные в consumer-датасете"),
        ("2", "OpenAI — Enterprise Privacy / data usage (2025)", "openai_terms",
         "consumer = обучение по умолчанию; API с мар 2023 не обучается на данных"),
        ("3", "EU AI Act — Reg. (EU) 2024/1689", "eu_ai_act",
         "штрафы: стандарт до 15M€/3%, верх до 35M€/7%"),
    ],
    "s24": [
        ("1", "Huang et al. — Survey on Hallucination in LLMs (2023)", "huang",
         "галлюцинация = уверенное порождение неверного, неотличимого от верного"),
        ("2", "Ji et al. — Survey of Hallucination in NLG (2023)", "ji",
         "таксономия галлюцинаций в генерации"),
        ("3", "Vectara — HHEM Leaderboard", "vectara",
         "диапазон <1% (суммаризация) → 10–15% (reasoning) — цифра зависит от задачи", True),
        ("4", "CybSafe / NCA — Oh Behave! (2024–25)", "cybsafe",
         "n=7000/7 стран: ~38% делятся конфиденциальным без ведома работодателя"),
    ],
    "s25": [
        ("1", "OpenAI — Sycophancy in GPT-4o postmortem (2025)", "sycophancy",
         "релиз 25 апр → откат 28 апр → разбор 29 апр; RLHF-переоценка приятных ответов"),
        ("2", "Pan et al. — Reward Misspecification (2022)", "reward_misspec",
         "reward hacking — общая природа: модель отражает данные/разметку"),
    ],
    "s26": [
        ("1", "Searle — Minds, Brains, and Programs (1980)", "searle",
         "Chinese Room: бенчмарк-эквивалентность ≠ понимание; narrow vs general"),
        ("2", "Bostrom — Superintelligence (2014)", "bostrom",
         "рамка долгосрочных AGI/ASI-сценариев для критического чтения прогнозов"),
    ],
}


# ============================================================
# Inline [N] injection into speaker notes — attach the marker right after the
# already-named textual attribution in each note (the notes name sources in
# words: «по бенчмарку Vectara», «Feng, McDonald, Zhang 2025» …). ORDER
# matters: longer/earlier phrases first. Each entry: (phrase, "[N]").
# The marker is inserted AFTER the first occurrence of phrase in the note body.
# ============================================================
NOTES_INLINE = {
    "s00b": [
        ("Gartner", "[1]"),
        ("пилот", "[2]"),
    ],
    "s01": [
        ("YOLOv8", "[1]"),
    ],
    "s06": [
        ("Russell & Norvig", "[1]"),
        ("ISO/IEC 22989", "[2]"),
        ("Mitchell", "[3]"),
        ("Сёрла", "[4]"),
    ],
    "s06a": [
        ("Уолтер Питтс", "[1]"),
    ],
    "s07": [
        ("Трансформер", "[1]"),
        ("AI Effect", "[2]"),
    ],
    "s08": [
        ("девятисот миллионов еженедельно активных пользователей", "[2]"),
        ("до сорока шести процентов кода написано AI", "[3] [VFY]"),
        ("Stack Overflow Developer Survey 2025", "[1]"),
        ("Grand View Research, 2026", "[4]"),
        ("Intellectual Analytics", "[5]"),
    ],
    "s09": [
        ("Mistral 7B превзошёл Llama-2 13B", "[1]"),
        ("MATH-500", "[2]"),
        ("SemiAnalysis", "[3]"),
        ("капитализация Nvidia упала на пятьсот восемьдесят девять миллиардов", "[4]"),
        ("OpenClaw", "[5]"),
        ("llama.cpp", "[6]"),
    ],
    "s11": [
        ("слоист", "[1]"),
        ("Agent", "[2]"),
    ],
    "s12": [
        ("классификаци", "[1]"),
    ],
    "s13": [
        ("контрол", "[1]"),
    ],
    "s15": [
        ("AlphaFold", "[2]"),
    ],
    "s16": [
        ("системный промпт", "[1]"),
    ],
    "s17": [
        ("ВЦИОМ", "[1]"),
        ("расширен до агента", "[2]"),
    ],
    "s18": [
        ("оркестратор", "[2]"),
    ],
    "s19": [
        ("последовательность вызов", "[1]"),
    ],
    "s19a": [
        ("Фэн, Макдональд и Чжан", "[1]"),
    ],
    "s20": [
        ("Google", "[1]"),
    ],
    "s21": [
        ("Два диагностических вопроса", "[1] [2] [3]"),
    ],
    "s22": [
        ("границ", "[1] [2] [3]"),
    ],
    "s23": [
        ("Samsung", "[1]"),
        ("Enterprise", "[2]"),
        ("EU AI Act", "[3]"),
    ],
    "s24": [
        ("Vectara Hughes Hallucination Evaluation Model", "[3]"),
        ("CybSafe", "[4]"),
    ],
    "s25": [
        ("апрель 2025", "[1]"),
        ("разметк", "[2]"),
    ],
    "s26": [
        ("narrow AI", "[1]"),
        ("Anthropic", "[2]"),
    ],
}


def _resolve_refs(sid):
    out = []
    for entry in SLIDE_REFS.get(sid, []):
        num, name, urlkey, gloss = entry[0], entry[1], entry[2], entry[3]
        volatile = len(entry) > 4 and entry[4]
        out.append((num, name, URLS.get(urlkey, ""), gloss, volatile))
    return out


def ref_list(slide, entries, *, y=6.95, x=0.55, w=12.25, h=0.50,
             size=8.5, color=LIGHT, line_spacing=1.02):
    """Bottom numbered CLICKABLE source list. entries: (num, name, url).
    Renders «[N] name» where name is a hyperlink. Muted, italic, small."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    for i, (num, name, url) in enumerate(entries):
        rm = p.add_run()
        rm.text = f"[{num}] "
        rm.font.name = FONT_BODY; rm.font.size = Pt(size)
        rm.font.bold = True; rm.font.italic = True
        rm.font.color.rgb = MID
        rn = p.add_run()
        rn.text = name
        rn.font.name = FONT_BODY; rn.font.size = Pt(size)
        rn.font.italic = True
        rn.font.color.rgb = color
        if url:
            try:
                rn.hyperlink.address = url
            except Exception:
                pass
        if i < len(entries) - 1:
            rs = p.add_run()
            rs.text = "   ·   "
            rs.font.name = FONT_BODY; rs.font.size = Pt(size)
            rs.font.italic = True
            rs.font.color.rgb = color
    return tb


def refs_of_slide(slide, sid, *, y=None, size=8.5):
    """Bottom clickable [N] list for a display slide, sourced from SLIDE_REFS.
    Rendered in a uniform footer band (y≈7.10) opened by nudging each slide's
    bottom callout/takeaway up (see build_sNN). Skips if no registry entry."""
    resolved = _resolve_refs(sid)
    if not resolved:
        return None
    entries = [(num, name, url) for (num, name, url, gloss, vol) in resolved]
    yy = y if y is not None else 7.06
    sz = size if len(entries) <= 4 else 7.6
    return ref_list(slide, entries, y=yy, size=sz, h=0.32)


def page_number(slide, n, total=None, *, color=SLATE):
    """Small muted page-number stamp in the bottom-right corner («N / TOTAL»).
    Applied to every slide by the assembler so all slides carry it."""
    txt = f"{n} / {total}" if total else str(n)
    tb = slide.shapes.add_textbox(Inches(12.33), Inches(7.16), Inches(0.95),
                                  Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = txt
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = color
    return tb


def notes_sources_block(sid):
    """Build the «Источники:» text block for the speaker notes: numbered [N] +
    FULL URL + one gloss phrase; volatile → [VFY-day-of]. "" if no entry."""
    resolved = _resolve_refs(sid)
    if not resolved:
        return ""
    lines = ["Источники:"]
    for (num, name, url, gloss, vol) in resolved:
        vfy = " [VFY-day-of]" if vol else ""
        lines.append(f"[{num}] {name} — {gloss}. {url}{vfy}")
    return "\n".join(lines)


def _inject_inline_refs(sid, body):
    """Insert inline [N] markers after the named textual attributions in the
    note body (NOTES_INLINE map). Idempotent-ish: inserts after the FIRST
    occurrence of each phrase that does not already carry the marker."""
    for phrase, marker in NOTES_INLINE.get(sid, []):
        idx = body.find(phrase)
        if idx < 0:
            continue
        end = idx + len(phrase)
        # skip if the marker already immediately follows
        if body[end:end + len(marker) + 1].strip().startswith(marker):
            continue
        body = body[:end] + " " + marker + body[end:]
    return body


def notes_with_sources(slide, sid):
    """Write speaker notes (paragraph-formatted) with inline [N] injected at
    named attributions AND the «Источники:» block appended. Single call
    replaces speaker_notes(slide, load_notes(sid))."""
    body = _inject_inline_refs(sid, load_notes(sid))
    block = notes_sources_block(sid)
    text = f"{body}\n\n{block}" if block else body
    speaker_notes(slide, text)


# ============================================================
# Roadmap bar (used by section dividers)
# ============================================================
def roadmap_bar(slide, here_idx):
    """Render a 5-section roadmap bar at bottom of slide.
    here_idx: 0=section 0 (open), 1=раздел 1, 2=раздел 2, 3=раздел 3, 4=раздел 4, 5=раздел 5."""
    # 6 cells (0..5) over 12.3 width, with 0.05 gaps
    bar_y = 6.55
    bar_h = 0.4
    n_cells = 6
    total_w = 12.3
    gap = 0.06
    cell_w = (total_w - gap * (n_cells - 1)) / n_cells
    start_x = 0.55
    labels = [
        "0  Введение",
        "1  AI",
        "2  Сейчас",
        "3  Способы",
        "4  Границы",
        "5  Итог",
    ]
    for i, label in enumerate(labels):
        x = start_x + i * (cell_w + gap)
        is_here = (i == here_idx)
        fill = GOLD if is_here else SOFT_GREY
        text_color = DEEP if is_here else SLATE
        filled_rect(slide, x, bar_y, cell_w, bar_h, fill, radius=True, radius_adj=0.30)
        text_box(slide, x=x, y=bar_y + 0.08, w=cell_w, h=bar_h - 0.16,
                 text=label, size=11, bold=is_here, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Fix-19: removed "Вы здесь — раздел N из 5" text. Color highlight (GOLD cell) is sufficient.


# ============================================================
# Unified navigation template (Fix-17 — single pattern across nav slides)
#
# Used by: s02a (overview, here=0), s10 (zoom-in, here=3),
#          s22 (zoom-in, here=4), s27 (zoom-in, here=5).
#
# Two states:
#   - overview (here_idx=0): all 6 cards equal weight, current=0 highlighted gold border.
#   - zoom-in (here_idx>=1): current card gold-FILLED + white text + frame phrase under title.
#
# Rationale: each nav slide shows ALL 5 sections (0..5 cards) so student can read
# the whole map and instantly see "where we are now". Identical visual structure
# across all dividers = predictable navigation, no re-orientation cost.
# (s14 deep-dive divider was deleted under Fix-17 — duplicated s10 framing.)
# ============================================================
NAV_SECTIONS = [
    # (num, title, short description). Used in `nav_slide` (s10/s22/s27).
    # issue #153 consistency fix: «и опросы» / «задание» removed — poll
    # moved to seminar 1 (fix #1) and homework callout removed from s28
    # (fix #18); these labels would otherwise contradict those removals.
    ("0", "Введение",                   "Демо · инструктор ·\nцентральный вопрос"),
    ("1", "Что такое AI",               "Определения,\nистория, перелом"),
    ("2", "Где мы\nсейчас",             "Цифры рынка\n2022–2026"),
    ("3", "Четыре способа\nреализации", "Модель · чат ·\nагент · приложение"),
    ("4", "Границы\nи безопасность",    "Что AI ломает\nи где не работает"),
    ("5", "Заключение",                 "Резюме ·\nкарта семестра"),
]


def nav_slide(slide, here_idx, title, frame_phrase=None, sub_marker=None):
    """Unified navigation slide layout.

    Parameters:
      here_idx     — 0 for overview state (s02a), 1..5 for zoom-in state.
                     - overview: cards equal-weight, gold border on card 0.
                     - zoom-in : current card gold-FILLED, white text inside.
      title        — slide title (centered, top).
      frame_phrase — optional 1-line italic frame under title (recommended for zoom-in).
      sub_marker   — DEPRECATED per Fix-19. Kept for backward-compat signature; ignored.
                     Color highlight on the active card is sufficient.
    """
    set_slide_bg(slide, SURFACE)
    # Title at top
    text_box(slide, x=0.55, y=0.45, w=12.25, h=0.95, text=title,
             size=30, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    if frame_phrase:
        text_box(slide, x=0.55, y=1.45, w=12.25, h=0.55,
                 text=frame_phrase, size=18, italic=True, color=MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.20)

    # 6 cards horizontal — same layout regardless of state for visual consistency
    card_y = 2.55
    card_w = 1.95
    card_h = 3.2
    gap = 0.15
    start_x = (SLIDE_W_IN - (card_w * 6 + gap * 5)) / 2.0

    for i, (num, sec_title, desc) in enumerate(NAV_SECTIONS):
        x = start_x + i * (card_w + gap)
        is_here = (i == here_idx)

        if is_here:
            # Highlighted card: gold-FILLED background, white text inside
            ocean_box(slide, x, card_y, card_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=2.5)
            num_color = WHITE
            title_color = WHITE
            desc_color = WHITE
        else:
            # Normal card: white fill, light stroke, normal colors
            ocean_box(slide, x, card_y, card_w, card_h,
                      fill=WHITE, stroke=LIGHT, stroke_pt=1.2)
            num_color = LIGHT if i < 2 else (MID if i < 4 else DEEP)
            title_color = DEEP
            desc_color = SLATE

        # Number — big at top
        text_box(slide, x=x, y=card_y + 0.30, w=card_w, h=0.85, text=num,
                 size=44, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        # Section title — middle
        text_box(slide, x=x + 0.08, y=card_y + 1.30, w=card_w - 0.16, h=1.05,
                 text=sec_title, size=14, bold=True, color=title_color,
                 align=PP_ALIGN.CENTER, line_spacing=1.20)
        # Description — bottom
        text_box(slide, x=x + 0.08, y=card_y + 2.30, w=card_w - 0.16, h=0.85,
                 text=desc, size=11, italic=True, color=desc_color,
                 align=PP_ALIGN.CENTER, line_spacing=1.25)

    # Fix-19: sub_marker (e.g. «Вы здесь») intentionally NOT rendered.
    # Color highlight (gold-filled card) is the only navigation indicator.
    _ = sub_marker  # parameter kept for backward-compat signature


# ============================================================
# Slide builders
# ============================================================
def build_s01(p):
    s = blank(p)
    text_box(s, x=0.55, y=0.55, w=5.9, h=2.4,
             text="Идентификация людей в реальном времени — на ноутбуке, без интернета, с 2023 года.",
             size=26, bold=True, color=DEEP, line_spacing=1.20)
    text_box(s, x=0.55, y=3.15, w=5.9, h=1.4,
             text="Narrow AI — модель решает одну задачу (обнаружение людей в кадре) и больше ничего.",
             size=15, italic=True, color=MID, line_spacing=1.3)
    # Bottom caption with mixed runs
    text_runs(s, 0.55, 5.5, 5.9, 1.0, [
        {"text": "На экране — ", "size": 15, "color": DEEP},
        {"text": "YOLOv8", "size": 15, "color": MID, "bold": True},
        {"text": " на CPU ноутбука: ", "size": 15, "color": DEEP},
        {"text": "~30 fps", "size": 15, "color": GOLD, "bold": True},
        {"text": ".", "size": 15, "color": DEEP},
        {"newpara": True, "text": "Без интернета", "size": 15, "color": TEAL, "bold": True},
        {"text": "  ·  обучена в 2023.", "size": 15, "color": DEEP},
        {"text": " [2]", "size": 15, "color": DEEP},
    ], line_spacing=1.35)
    # Right column — Ocean rounded box framing the YOLO mock screenshot
    box_x, box_y, box_w, box_h = 6.55, 0.55, 6.3, 4.4
    ocean_box(s, box_x, box_y, box_w, box_h)
    pad = 0.18
    img_w = box_w - 2 * pad
    img_h = img_w * 720.0 / 1280.0
    img_x = box_x + pad
    img_y = box_y + (box_h - img_h) / 2.0
    add_image(s, ASSETS / "illustrations/s01-yolo-mock.png", img_x, img_y, img_w, img_h)
    text_box(s, x=box_x, y=box_y + box_h + 0.05, w=box_w, h=0.4,
             text="Кадр модели в момент демо: 2 человека в боксах. YOLOv8 (Ultralytics, 2023). [1]",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    refs_of_slide(s, "s01")
    notes_with_sources(s, "s01")


def build_s00a(p):
    """Welcome hero — issue #153 fix #2. Short greeting before cover.

    Hero/title composition in Ocean palette, tinted bg, large typography —
    similar spirit to cover but WITHOUT the Ocean rounded box motif (motif
    is content-slide only) and without the decorative «01».
    """
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.9, y=2.15, w=11.5, h=1.0,
             text="Добро пожаловать на курс",
             size=30, bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.15)
    text_box(s, x=0.9, y=2.85, w=11.5, h=1.6,
             text="«Отраслевое применение систем\nискусственного интеллекта»",
             size=36, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.25)
    filled_rect(s, SLIDE_W_IN / 2 - 0.5, 4.55, 1.0, 0.06, fill=GOLD)
    text_box(s, x=1.4, y=4.85, w=10.5, h=0.7,
             text="17 лекций о том, где AI работает в индустриях — и где нет.",
             size=20, italic=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.25)
    speaker_notes(s, load_notes("s00a"))


def build_s00b(p):
    """Course hook — issue #153 fix #2. Funnel + central question, moved
    before cover from old s05b (reworded role: hook for engagement, not
    "course frame after instructor").
    """
    s = blank(p)
    slide_title(s, "Главный вопрос курса — не «можно ли AI?», а «нужно ли и где?».", size=24)
    # Left: funnel — same visual as old s05b (Fix-6 sizing preserved)
    fun_x, fun_y, fun_w = 0.55, 2.05, 5.5
    blk_h = 1.05
    blk_gap = 0.10
    filled_rect(s, fun_x, fun_y, fun_w, blk_h, LIGHT, radius=True, radius_adj=0.08)
    text_box(s, x=fun_x, y=fun_y + 0.20, w=fun_w, h=blk_h - 0.40,
             text="100% AI-пилотов запускаются",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)
    mid_w = fun_w * 0.75
    mid_x = fun_x + (fun_w - mid_w) / 2.0
    mid_y = fun_y + blk_h + blk_gap
    filled_rect(s, mid_x, mid_y, mid_w, blk_h, MID, radius=True, radius_adj=0.10)
    text_box(s, x=mid_x, y=mid_y + 0.20, w=mid_w, h=blk_h - 0.40,
             text="−90% откатываются",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bot_w = fun_w * 0.45
    bot_x = fun_x + (fun_w - bot_w) / 2.0
    bot_y = mid_y + blk_h + blk_gap
    filled_rect(s, bot_x, bot_y, bot_w, blk_h, GOLD, radius=True, radius_adj=0.10)
    text_box(s, x=bot_x, y=bot_y + 0.14, w=bot_w, h=blk_h - 0.28,
             text="10% доходят\nдо прода",
             size=18, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.10)
    text_box(s, x=fun_x, y=bot_y + blk_h + 0.15, w=fun_w, h=0.5,
             text="Иллюстрация принципа (Gartner, McKinsey подтверждают похожие цифры). [1]",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Right: takeaway + central question
    right_x = 6.3
    right_w = 6.5
    ocean_box(s, right_x, fun_y - 0.3, right_w, 5.0)
    text_box(s, x=right_x + 0.3, y=fun_y, w=right_w - 0.6, h=0.45,
             text="Главная мысль курса",
             size=14, bold=True, color=TEAL)
    text_box(s, x=right_x + 0.3, y=fun_y + 0.5, w=right_w - 0.6, h=1.6,
             text="Завтра — почти везде.\nСегодня — почти никто.\nКурс — про этот разрыв. [2]",
             size=22, bold=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=right_x + 0.3, y=fun_y + 2.4, w=right_w - 0.6, h=0.45,
             text="Центральный вопрос курса",
             size=14, bold=True, color=DEEP)
    text_box(s, x=right_x + 0.3, y=fun_y + 2.85, w=right_w - 0.6, h=1.6,
             text="Где AI работает,\nгде — нет,\nи как это понять?",
             size=24, bold=True, color=DEEP, line_spacing=1.25)
    refs_of_slide(s, "s00b")
    notes_with_sources(s, "s00b")


def build_s02(p):
    """Cover — distinct: tinted bg, decorative «01», 60pt title.

    issue #153 QA fix #3 (P0, presentation-critic): "1" was fully invisible —
    root cause was NOT only z-order. At 320pt bold, the LibreOffice-rendered
    fallback font (Arial unavailable → DejaVu Sans Bold substitute) makes
    "01" ≈6.2" wide, wider than the 5.3" text_box — word_wrap=True then
    wraps "1" onto a second line, which falls outside the box's usable
    height (single 320pt line ≈4.44" already nearly fills the 4.7" box) and
    is invisible (clipped), NOT merely covered by the hero image. Fix: (a)
    reduce font 320pt→230pt so "01" renders on one line with margin inside
    the box width (verified via PIL font-metrics: "01" @230pt ≈4.44" wide,
    fits inside 5.3"); (b) ALSO draw the hero image BEFORE the numeral
    text_box (was: image added after, on top) as defense-in-depth so any
    future numeral/illustration overlap resolves with text on top.
    """
    s = blank(p)
    set_slide_bg(s, SURFACE)
    if (ASSETS / "illustrations/hero-cover-light.png").exists():
        add_image(s, ASSETS / "illustrations/hero-cover-light.png",
                  x=8.0, y=0.9, w=5.0, h=5.0)
    text_box(s, x=8.0, y=2.9, w=5.3, h=3.6, text="01",
             size=230, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.7, y=1.0, w=6.5, h=0.55, text="ЛЕКЦИЯ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.72, 1.55, 0.7, 0.05, fill=TEAL)
    text_box(s, x=0.7, y=2.0, w=8.5, h=2.4, text="Что такое AI?\nИстория, классификация,\nобщие понятия",
             size=42, bold=True, color=DEEP, line_spacing=1.10, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.7, 5.45, 0.05, 0.55, fill=TEAL)
    text_box(s, x=0.95, y=5.45, w=8.0, h=0.6,
             text="Карта применений AI: где работает, где — нет.",
             size=22, color=MID, italic=False, align=PP_ALIGN.LEFT, line_spacing=1.25)
    speaker_notes(s, load_notes("s02"))


def build_s02a(p):
    """Lecture map — issue #153 fix #3: REDESIGN as horizontal timeline.

    Was: 6-card equal-weight grid (unified nav template, same skeleton as
    s10/s22/s27 dividers). Now: horizontal timeline with colored zone blocks,
    visually matching the s29 course-roadmap redesign pattern (modular
    color blocks instead of equal cards) — while keeping THIS slide's
    content scoped to the lecture (5 sections), not the course.
    Content updated: «Открытие и опросы» → «Открытие» (poll removed, fix #1);
    section 3 stays the biggest zone (widest block) reflecting its size.
    """
    s = blank(p)
    slide_title(s, "План лекции", size=28, align=PP_ALIGN.CENTER)
    sections = [
        ("0", "Введение", "Демо · инструктор ·\nцентральный вопрос", TEAL, 1.0, True),
        ("1", "Что такое AI", "Определения,\nистория, перелом", LIGHT, 1.0, False),
        ("2", "Где мы сейчас", "Цифры рынка\n2022–2026", MID, 1.0, False),
        ("3", "4 способа реализации", "Модель · чат ·\nагент · приложение", DEEP, 1.6, False),
        ("4", "Границы\nи безопасность", "Что AI ломает\nи где не работает", MID, 1.2, False),
        ("5", "Заключение", "Резюме · карта\nсеместра", LIGHT, 1.0, False),
    ]
    mod_y = 2.0
    mod_h = 4.3
    bar_x = 0.55
    bar_w = SLIDE_W_IN - 2 * 0.55
    total_units = sum(u for *_, u, _ in sections)
    unit_w = bar_w / total_units
    cur_x = bar_x
    for num, title_txt, desc, color, units, is_now in sections:
        m_w = units * unit_w
        fill_color = GOLD if is_now else color
        text_color = DEEP if is_now else WHITE
        ocean_box(s, cur_x, mod_y, m_w - 0.06, mod_h, fill=WHITE, stroke=color, stroke_pt=2.0)
        # issue #155 QA fix P2-9 (s02a/s27 nav-card sliver): the colored
        # header strip's rounded corners (fixed radius_adj=0.10 on a 1.0"-tall
        # shape → ~0.05" absolute radius) were noticeably tighter than the
        # outer ocean_box card's corners (12pt ≈ 0.167" absolute radius, per
        # ocean_box's own radius_pt formula) — the header's straighter corner
        # sat fully inside the card's rounder corner, so a thin sliver of the
        # card's own stroke colour showed through the gap at the top corners
        # (visible on EVERY card, not just the edges — misread on first look
        # as an edge-only artifact). Fix: give the header strip the same
        # absolute ~12pt corner radius as ocean_box uses, computed against
        # its own (shorter) height so both shapes round by the same physical
        # amount.
        strip_h = 1.0
        strip_radius_adj = max(0.04, min(0.35, (12.0 / 72.0) / (strip_h / 2.0)))
        filled_rect(s, cur_x, mod_y, m_w - 0.06, strip_h, fill_color, radius=True, radius_adj=strip_radius_adj)
        text_box(s, x=cur_x, y=mod_y + 0.08, w=m_w - 0.06, h=0.55, text=num,
                 size=30, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        text_box(s, x=cur_x + 0.10, y=mod_y + 1.15, w=m_w - 0.26, h=0.75, text=title_txt,
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
        text_box(s, x=cur_x + 0.10, y=mod_y + 2.05, w=m_w - 0.26, h=1.0, text=desc,
                 size=10.5, italic=True, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.25)
        cur_x += m_w
    speaker_notes(s, load_notes("s02a"))


def build_s05a(p):
    """Instructor card — issue #155 fix #174: full rebuild on sem-01 s02 reference.

    Layout: left vertical strip (~28% width, SURFACE fill) with portrait
    photo + full name + divider + 2 contact rows (Telegram/Email badges).
    Right area (~72%, white) with specialization headline + 3 fact cards
    (experience numbers / expertise / company badges) in a 2-tier layout,
    matching library/seminars/sem-01/slides/s02-instructor-bio.md pattern.
    """
    s = blank(p)
    # ---- Left strip ----
    strip_w = 3.75
    filled_rect(s, 0, 0, strip_w, SLIDE_H_IN, SURFACE)
    filled_rect(s, strip_w - 0.02, 0, 0.02, SLIDE_H_IN, LIGHT)
    photo_x, photo_y, photo_w, photo_h = 0.55, 0.55, 2.65, 3.53
    add_image(s, ASSETS / "instructor-photo-crop.png",
              x=photo_x, y=photo_y, w=photo_w, h=photo_h)
    # thin frame around photo
    frame = filled_rect(s, photo_x, photo_y, photo_w, photo_h, WHITE, stroke=LIGHT, stroke_pt=1.5)
    frame.fill.background()
    name_y = photo_y + photo_h + 0.25
    text_box(s, x=0.35, y=name_y, w=strip_w - 0.6, h=0.85,
             text="Левко Максим\nНиколаевич", size=20, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.15)
    # divider line
    div_y = name_y + 0.95
    filled_rect(s, 0.55, div_y, strip_w - 1.1, 0.02, COVER_OUTLINE)
    # contact rows
    contacts = [
        ("lucide-send-blue.png", "Telegram", "@Maxim_Levko"),
        ("lucide-mail-blue.png", "Email", "Levko.maxim@gmail.com"),
    ]
    cy = div_y + 0.30
    for icon, label, value in contacts:
        badge = filled_rect(s, 0.55, cy, 0.5, 0.5, WHITE, stroke=LIGHT, stroke_pt=1.2,
                            radius=True, radius_adj=0.5)
        add_image(s, ASSETS / "icons" / icon, x=0.55 + 0.11, y=cy + 0.11, w=0.28, h=0.28)
        text_box(s, x=1.20, y=cy - 0.02, w=strip_w - 1.35, h=0.28,
                 text=label, size=11, bold=True, color=TEAL)
        text_box(s, x=1.20, y=cy + 0.24, w=strip_w - 1.35, h=0.30,
                 text=value, size=11.5, color=DEEP)
        cy += 0.72
    # ---- Right area ----
    rx = strip_w + 0.45
    rw = SLIDE_W_IN - rx - 0.55
    slide_title(s, "Кто я и почему мне это важно.", size=26, x=rx, w=rw, y=0.45, h=0.85)
    text_box(s, x=rx, y=1.30, w=rw, h=0.65,
             text="Архитектор, технический и продуктовый лидер создания\nи внедрения информационных систем",
             size=15, bold=True, color=MID, line_spacing=1.20)
    # Card tier 1: experience stat + expertise (2 cards side by side)
    card_y1 = 2.35
    card_h1 = 2.15
    gap = 0.25
    card_w1 = (rw - gap) / 2
    # Card A — experience numbers
    ocean_box(s, rx, card_y1, card_w1, card_h1)
    add_image(s, ASSETS / "icons" / "lucide-briefcase-blue.png",
              x=rx + 0.30, y=card_y1 + 0.32, w=0.55, h=0.55)
    text_runs(s, x=rx + 0.30, y=card_y1 + 1.02, w=card_w1 - 0.6, h=0.40,
              runs=[
                  {"text": "20+", "size": 16, "bold": True, "color": GOLD},
                  {"text": " лет опыта в ИТ", "size": 16, "bold": True, "color": DEEP},
              ])
    text_box(s, x=rx + 0.30, y=card_y1 + 1.48, w=card_w1 - 0.6, h=0.55,
             text="10+ завершённых проектов\nпод руководством", size=11.5, italic=True, color=SLATE, line_spacing=1.25)
    # Card B — expertise
    bx = rx + card_w1 + gap
    ocean_box(s, bx, card_y1, card_w1, card_h1)
    add_image(s, ASSETS / "icons" / "lucide-layers-blue.png",
              x=bx + 0.30, y=card_y1 + 0.32, w=0.55, h=0.55)
    text_box(s, x=bx + 0.30, y=card_y1 + 1.02, w=card_w1 - 0.6, h=0.35,
             text="Экспертиза", size=16, bold=True, color=DEEP)
    text_box(s, x=bx + 0.30, y=card_y1 + 1.42, w=card_w1 - 0.6, h=0.65,
             text="Системный анализ · Проектирование систем · Управление данными · Автоматизация бизнеса · Управление продуктами",
             size=10, italic=True, color=SLATE, line_spacing=1.25)
    # Card tier 2: company badges (wide card, generic icons, no logos)
    card_y2 = card_y1 + card_h1 + 0.30
    card_h2 = 1.55
    ocean_box(s, rx, card_y2, rw, card_h2)
    text_box(s, x=rx + 0.30, y=card_y2 + 0.22, w=rw - 0.6, h=0.35,
             text="Консалтинг и inhouse", size=16, bold=True, color=DEEP)
    companies = ["Yandex", "МТС", "Магнит", "Сибур"]
    pill_h = 0.62
    pill_gap = 0.20
    pill_w = (rw - 0.6 - 3 * pill_gap) / 4
    for i, comp in enumerate(companies):
        px = rx + 0.30 + i * (pill_w + pill_gap)
        py = card_y2 + 0.78
        filled_rect(s, px, py, pill_w, pill_h, TEAL_TINT, stroke=TEAL, stroke_pt=1.2,
                   radius=True, radius_adj=0.5)
        add_image(s, ASSETS / "icons" / "lucide-building-2-blue.png",
                  x=px + 0.14, y=py + pill_h/2 - 0.16, w=0.32, h=0.32)
        text_box(s, x=px + 0.52, y=py, w=pill_w - 0.60, h=pill_h,
                 text=comp, size=12.5, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s05a"))


def build_s05c(p):
    """Section 1 divider — issue #155 fix #177 (section-divider audit).

    Zoom-in state of the unified nav template (same pattern as s10/s22/s27):
    6-card grid, card 1 gold-FILLED with white text. Added because the
    divider audit found NO divider before раздел 1 «Что такое AI» (lecture
    jumped straight from s05a instructor card into s06 content).
    """
    s = blank(p)
    nav_slide(s, here_idx=1,
              title="Раздел 1 · Что такое AI",
              frame_phrase="Определения, история, классификация.")
    speaker_notes(s, load_notes("s05c"))


def build_s06(p):
    """Multiple definitions of AI — 4 approaches grid (full definitions, not labels) + AI Effect.

    Fix-13 (2026-05-13): cards now contain the actual definitions (~15-20 words each)
    instead of just approach names. Body 14pt (was 12pt), source 11pt italic (was 10pt),
    cell_h 2.40" (was 1.95") to fit full text. Grid moved up; callout moved down.
    """
    s = blank(p)
    slide_title(s, "Определений AI много — потому что AI это движущаяся цель.", size=26)
    cards = [
        ("Russell & Norvig (AIMA, 2021)",
         "«AI = система, мыслящая как человек, мыслящая рационально, действующая как человек или действующая рационально (4 квадранта по 2 осям).»",
         "Russell & Norvig, AIMA, 4-е изд., 2021 [1]",
         MID),
        ("ISO/IEC 22989:2022",
         "«AI-система — это engineered system, которая генерирует выходы (рекомендации, прогнозы, решения) для целей, заданных человеком.»",
         "Международный стандарт ISO/IEC 22989:2022 — опора EU AI Act [2]",
         LIGHT),
        ("Через обучение (Mitchell, 1997)",
         "«Программа улучшается с опытом E на задаче T по метрике P. Если поведение возникает из обученной модели — это AI.»",
         "Mitchell, Machine Learning, 1997 [3]",
         MID),
        ("Через бенчмарки и AGI",
         "«AI = то, что проходит тест Тьюринга или решает бенчмарк на уровне человека.» Возражение Сёрла: поведение ≠ понимание.",
         "Тьюринг 1950 / Searle 1980 — Chinese Room [4]",
         LIGHT),
    ]
    grid_x = 0.55
    grid_y = 1.62
    cell_w = 6.05
    cell_h = 2.40
    cell_gap = 0.15
    for i, (head, body, src, color) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = grid_x + col * (cell_w + cell_gap)
        y = grid_y + row * (cell_h + cell_gap)
        ocean_box(s, x, y, cell_w, cell_h, stroke=color)
        # Header (approach name)
        text_box(s, x=x + 0.25, y=y + 0.18, w=cell_w - 0.5, h=0.40,
                 text=head, size=15, bold=True, color=color)
        # Body — full definition, 14pt for projector readability
        text_box(s, x=x + 0.25, y=y + 0.60, w=cell_w - 0.5, h=1.40,
                 text=body, size=14, color=DEEP, line_spacing=1.22)
        # Source citation
        text_box(s, x=x + 0.25, y=y + cell_h - 0.40, w=cell_w - 0.5, h=0.32,
                 text=src, size=11, italic=True, color=SLATE)
    # AI Effect callout at bottom (gold accent, ≥1×/slide rule)
    gold_callout(s, 0.55, 6.42, 12.25, 0.55,
                 "AI Effect (Tesler):  «AI is whatever hasn't been done yet».  Как только техника начинает работать, её перестают называть AI.",
                 size=13)
    refs_of_slide(s, "s06")
    notes_with_sources(s, "s06")


def build_s06a(p):
    """New fact-bridge slide — issue #153 fix #4.

    McCulloch-Pitts 1943 vs Dartmouth 1956 — 13-year gap. Short, compact,
    single fact anchor between s06 (definitions) and s07 (70-year timeline).

    issue #155 comment #175: added a real photo sidebar on the right —
    Walter Pitts with Jerome Lettvin (1959 frog-experiment photo, Wikimedia
    Commons, CC BY-SA 3.0). Owner-approved substitute after the originally
    requested "McCulloch+Pitts 1943 joint photo" URL turned out to be a
    mislabeled 1963 Nobel Prize cover with unrelated scientists — no
    clean-license solo portrait of either man exists. This is a GROUP photo
    from 1959, NOT 1943, so the caption explicitly disclaims the year
    mismatch (see slide .md Visual section + iteration-log-issue155.md
    "s06a photo insert" for the full acquisition trail).

    Layout change to make room: the two year-anchor boxes shrank from
    3.7"→2.95" wide and moved into the left ~65% of the canvas (freeing a
    ~3.3" right-hand column for the photo). The gold callout + closing
    takeaway line stay full-width at the bottom, unchanged in content.
    """
    s = blank(p)
    slide_title(s, "Идея нейросети старше самого термина «искусственный интеллект» на 13 лет.", size=23,
                w=12.3)
    # ---- Right sidebar: real 1959 photo (Pitts + Lettvin) with caption ----
    # Source jpg pre-cropped (see notes/mcp-limitations.md-adjacent acquisition
    # log) to trim excess dark background top/bottom: 490x569 px, ratio 0.861.
    photo_h = 2.30
    photo_w = photo_h * (490 / 569)
    photo_x = SLIDE_W_IN - 0.55 - photo_w
    photo_y = 2.35
    add_image(s, ROOT / "assets/images/lettvin-pitts-1959-crop.jpg",
              x=photo_x, y=photo_y, w=photo_w, h=photo_h)
    frame = filled_rect(s, photo_x, photo_y, photo_w, photo_h, WHITE, stroke=LIGHT, stroke_pt=1.5)
    frame.fill.background()
    text_box(s, x=photo_x - 0.15, y=photo_y + photo_h + 0.07, w=photo_w + 0.30, h=0.65,
             text="Питтс (справа) и Леттвин, 1959 — опыт с лягушкой в MIT "
                  "(это НЕ фото 1943 г.). Wikimedia Commons, CC BY-SA 3.0.",
             size=7.8, italic=True, color=SLATE, align=PP_ALIGN.LEFT, line_spacing=1.1)
    # ---- Left: two year anchors with a gold bridge between them ----
    # anchor_w solved so [left_x .. right anchor's right edge] ends exactly
    # gap_to_photo (0.65") before the photo's left edge — no leftover dead
    # whitespace between the two visual groups.
    anchor_y = 2.35
    anchor_h = 2.35
    left_x = 0.55
    gap_w = 1.3
    gap_to_photo = 0.55
    anchor_w = ((photo_x - gap_to_photo) - left_x - gap_w) / 2
    right_x = left_x + anchor_w + gap_w
    ocean_box(s, left_x, anchor_y, anchor_w, anchor_h, stroke=LIGHT)
    text_box(s, x=left_x, y=anchor_y + 0.28, w=anchor_w, h=0.9, text="1943",
             size=46, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_box(s, x=left_x + 0.24, y=anchor_y + 1.28, w=anchor_w - 0.48, h=0.95,
             text="Мак-Каллок и Питтс [1]:\nформальный нейрон\nкак логический элемент",
             size=11.5, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.22)
    ocean_box(s, right_x, anchor_y, anchor_w, anchor_h, stroke=MID)
    text_box(s, x=right_x, y=anchor_y + 0.28, w=anchor_w, h=0.9, text="1956",
             size=46, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(s, x=right_x + 0.24, y=anchor_y + 1.28, w=anchor_w - 0.48, h=0.95,
             text="Дартмутская конференция:\nтермин «artificial intelligence»",
             size=11.5, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.22)
    # Gold bridge arrow + "13 лет" label between the two anchors
    bridge_x = left_x + anchor_w + 0.15
    bridge_w = right_x - bridge_x - 0.15
    bridge_y = anchor_y + 0.55
    filled_rect(s, bridge_x, bridge_y, bridge_w, 0.20, GOLD, radius=True, radius_adj=0.5)
    text_box(s, x=bridge_x - 0.25, y=anchor_y - 0.05, w=bridge_w + 0.5, h=0.5,
             text="13 лет", size=20, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    gold_callout(s, 0.55, 5.30, 12.25, 0.95,
                 "Формальный нейрон не решал прикладных задач, но предвосхитил коннекционистскую традицию — линию мысли, из которой десятилетия спустя вырастут нейросети и Трансформер.",
                 size=13)
    text_box(s, x=0.55, y=6.48, w=12.25, h=0.5,
             text="Идея «нейросети» на теоретическом уровне старше самого термина «искусственный интеллект».",
             size=14, italic=True, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.25)
    refs_of_slide(s, "s06a")
    notes_with_sources(s, "s06a")


def build_s07(p):
    """70 years AI timeline — issue #155 Round-2 redesign (owner comment #176).

    Owner feedback on Fix-7 v2 (batch 1-4 baseline): left group-labels
    duplicated the year-range already shown on the timeline itself
    ("Открытия (1950 — 1980-е)" next to a line that already shows 1950 /
    1966 / 1980-е), and the 3 groups had no visual separation — just 3
    stacked lines on white.

    Redesign (5 iterations, Generate→Convert→Inspect→Fix, tried 3
    substantially different compositions — see iteration-log-issue155.md
    "s07 Round-2 redesign" for the full a/b/c/d trail):
    - Left text-column label REMOVED entirely — year-range dropped since
      years are already on the line; group name became a badge.
    - Each of the 3 groups now sits on its own full-width tinted panel
      (Ocean light/mid/deep tint backgrounds), with a rounded WHITE-on-color
      "tab" pill straddling the panel's top edge as the group title — this
      is the group's visual home instead of a separate left column.
    - WCAG fix found during iteration: gold-COLORED TEXT on a light tinted
      background measures ~1.6:1 contrast (fails AA) — gold in this deck
      only works as a FILL with dark text on top (~6.9:1), never as text
      color on a light bg. Pivot year "2017" text and pivot event label are
      now DEEP (not gold); the "wow" signal moved to a bigger gold pill +
      oval marker straddling the timeline (reads as a distinct badge/pin).
    - Turing label Russified ("Turing — Imitation Game" → "Тьюринг — тест
      на мышление") per owner request; Weizenbaum/Lighthill Russified to
      match the spelling already used in this slide's own speaker notes
      ("Вайценбаум", "Лайтхилла" — see s07-timeline-2017.md speaker notes).
      "«Attention Is All You Need»" kept verbatim — exact paper title in
      quotes, legitimate citation, not an anglicism.
    - All facts (dates, events, gold-2017 accent, Vaswani callout with all
      7 co-authors + 160K+ citations) unchanged — content contract preserved.
    """
    s = blank(p)
    slide_title(s, "70 лет AI: открытия, зимы, точка перелома 2017. [2] [3]", size=28)

    def tint(color_tuple, factor=0.90):
        r = int(color_tuple[0] + (255 - color_tuple[0]) * factor)
        g = int(color_tuple[1] + (255 - color_tuple[1]) * factor)
        b = int(color_tuple[2] + (255 - color_tuple[2]) * factor)
        return RGBColor(r, g, b)

    LIGHT_TINT_BG = tint((0x1C, 0x72, 0x93))
    MID_TINT_BG = tint((0x06, 0x5A, 0x82))
    DEEP_TINT_BG = tint((0x21, 0x29, 0x5C))

    groups = [
        ("Открытия", LIGHT, LIGHT_TINT_BG, None, [
            ("1950", "Тьюринг — тест на мышление"),
            ("1966", "ELIZA — Вайценбаум"),
            ("1980-е", "Экспертные системы"),
        ]),
        # issue #155 QA fix P2-10: the "zимы" takeaway ("ресурсы уходят, когда
        # обещания не сбываются") previously lived ONLY in speaker notes —
        # invisible on the slide itself. Added as a short inline caption next
        # to the group-name tab (not GOLD text per P1-1 WCAG rule — DEEP is
        # readable on both the white slide bg behind the tab row and the
        # MID_TINT_BG panel it sits just above).
        ("Зимы и прорывы", MID, MID_TINT_BG, "ресурсы уходят, когда обещания не сбываются", [
            ("1974", "1-я зима — доклад Лайтхилла"),
            ("1997", "Deep Blue — 200M поз/сек"),
            ("2012", "AlexNet — GPU + DL"),
        ]),
        ("Перелом и взрыв", DEEP, DEEP_TINT_BG, None, [
            ("2017", "«Attention Is All You Need»  ★"),
            ("2022", "ChatGPT — 1M за 5 дней"),
            ("2025-26", "DeepSeek R1, Claude Code"),
        ]),
    ]
    band_h = 1.48
    band_y_start = 1.72
    gap = 0.15
    panel_x, panel_w = 0.55, 12.25
    for gi, (gname, color, bg, caption, events) in enumerate(groups):
        band_y = band_y_start + gi * (band_h + gap)
        # Full-width tinted panel — the group's visual "home" (replaces the
        # old plain-white background + left text column).
        filled_rect(s, panel_x, band_y, panel_w, band_h, bg, radius=True, radius_adj=0.14)
        # Group-name tab, straddling the panel's top edge.
        tab_w = 0.42 + 0.155 * len(gname)
        filled_rect(s, panel_x + 0.35, band_y - 0.19, tab_w, 0.40, color,
                    radius=True, radius_adj=0.5)
        text_box(s, x=panel_x + 0.35, y=band_y - 0.19, w=tab_w, h=0.40, text=gname,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if caption:
            cap_x = panel_x + 0.35 + tab_w + 0.20
            text_box(s, x=cap_x, y=band_y - 0.19, w=panel_w - (cap_x - panel_x) - 0.30, h=0.40,
                     text=caption, size=10.5, italic=True, color=DEEP,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        line_y = band_y + band_h - 0.44
        line_x0, line_w = panel_x + 0.45, panel_w - 0.90
        filled_rect(s, line_x0, line_y, line_w, 0.08, color, radius=True, radius_adj=0.5)
        n = len(events)
        ev_w = line_w / n
        for ei, (year, label) in enumerate(events):
            ex = line_x0 + ei * ev_w
            is_pivot = "★" in label
            tick_x = ex + ev_w / 2 - 0.08
            if is_pivot:
                # issue #155 QA fix P2-9: oval marker was centred ON the
                # timeline (top line_y-0.13, bottom line_y+0.21) while the
                # gold pill below started at line_y-0.02 — the pill covered
                # the bottom ~2/3 of the oval, leaving only a thin
                # DEEP-stroked arc poking out above the pill (read as a
                # rendering glitch on inspection, not the intended
                # "badge/pin silhouette"). Raised the oval so its bottom
                # edge (line_y-0.40+0.34 = line_y-0.06) clears the pill's
                # top (line_y-0.02) with a small visible gap (~0.04").
                shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(tick_x - 0.11), Inches(line_y - 0.40),
                                         Inches(0.34), Inches(0.34))
                shp.fill.solid(); shp.fill.fore_color.rgb = GOLD
                shp.line.color.rgb = DEEP; shp.line.width = Pt(1.75)
                disable_shadow(shp)
            else:
                filled_rect(s, tick_x + 0.04, line_y - 0.07, 0.08, 0.22, color)
            # Event label — DEEP for both pivot and regular (WCAG fix: gold
            # text on light tint bg measured ~1.6:1, fails AA; see docstring).
            text_box(s, x=ex, y=band_y + 0.33, w=ev_w, h=0.42, text=label,
                     size=12.5, color=DEEP,
                     bold=is_pivot, align=PP_ALIGN.CENTER, line_spacing=1.05)
            if is_pivot:
                # Pivot year sits in its own gold pill (DEEP text on GOLD
                # fill = ~6.9:1 contrast) — bigger + shape-distinct so it
                # reads as the slide's single visual anchor.
                pill_w = 1.35
                filled_rect(s, ex + ev_w / 2 - pill_w / 2, band_y + band_h - 0.46,
                            pill_w, 0.44, GOLD, radius=True, radius_adj=0.5)
                text_box(s, x=ex, y=band_y + band_h - 0.45, w=ev_w, h=0.40,
                         text=year, size=22, bold=True,
                         color=DEEP, align=PP_ALIGN.CENTER)
            else:
                text_box(s, x=ex, y=band_y + band_h - 0.35, w=ev_w, h=0.32,
                         text=year, size=14.5, bold=True,
                         color=color, align=PP_ALIGN.CENTER)
    # Vaswani-2017 deep-dive callout — unchanged content (all 7 co-authors +
    # 160K+ citations), repositioned for the new panel geometry.
    # Bands occupy: 1.72 + 3×1.48 + 2×0.15 = 6.46. Callout starts at 6.66.
    gold_callout(s, 0.55, 6.26, 12.25, 0.72,
                 "★ 2017 — Vaswani и 7 соавторов (Shazeer, Parmar, Uszkoreit, Jones, Gomez, Kaiser, Polosukhin) "
                 "вводят механизм self-attention — основу всех современных LLM. "
                 "На май 2026 у статьи свыше 160 000 цитирований. [1]",
                 size=12.5)
    refs_of_slide(s, "s07")
    notes_with_sources(s, "s07")


def build_s07a(p):
    """Section 2 divider — issue #155 fix #177 (section-divider audit).

    Zoom-in state of the unified nav template (same pattern as s10/s22/s27):
    6-card grid, card 2 gold-FILLED with white text. Added because the
    divider audit found NO divider before раздел 2 «Где мы сейчас» (lecture
    jumped straight from s07 timeline into s08 scale-numbers content).
    """
    s = blank(p)
    nav_slide(s, here_idx=2,
              title="Раздел 2 · Где мы сейчас",
              frame_phrase="AI как инфраструктура, прорывы 2023–2026.")
    speaker_notes(s, load_notes("s07a"))


def build_s08(p):
    """Scale numbers — 4 metrics grid + counter-fact gold."""
    s = blank(p)
    slide_title(s, "AI стал инфраструктурой за 3 года: 900M пользователей, 51% разработчиков ежедневно, 46% Copilot-кода.", size=22)
    metrics = [
        ("900M", "WAU", "ChatGPT, февраль 2026 [2]", "OpenAI", MID, "lucide-users-2-blue.png"),
        ("51%", "ежедневно", "Stack Overflow Dev Survey 2025 [1]", "n=49k+, 177 стран", LIGHT, "lucide-code-blue.png"),
        ("46%", "кода у Copilot", "GitHub Octoverse 2025 [3]", "Java — 61%", MID, "lucide-github.png"),
        ("$390.9B→$539.5B", "AI-рынок 2025→2026", "Grand View Research, 2026 [4]", "Statista (software-only): ~$244–260B", LIGHT, "lucide-dollar-sign-blue.png"),
    ]
    grid_y = 2.0
    cell_w = 6.05
    cell_h = 1.85
    grid_x = 0.55
    cell_gap = 0.15
    for i, (big, label, src1, src2, color, icon) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = grid_x + col * (cell_w + cell_gap)
        y = grid_y + row * (cell_h + cell_gap)
        ocean_box(s, x, y, cell_w, cell_h, stroke=color)
        # Icon top right
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + cell_w - 0.85, y=y + 0.25, w=0.55, h=0.55)
        # Big number — smaller font for the longer market-size string (fix #6)
        big_size = 30 if len(big) > 10 else 44
        text_box(s, x=x + 0.30, y=y + 0.20, w=cell_w - 1.4, h=0.85,
                 text=big, size=big_size, bold=True, color=color, line_spacing=1.0)
        text_box(s, x=x + 0.30, y=y + 0.95, w=cell_w - 0.5, h=0.4,
                 text=label, size=15, bold=True, color=DEEP)
        text_box(s, x=x + 0.30, y=y + cell_h - 0.55, w=cell_w - 0.5, h=0.32,
                 text=src1, size=10, italic=True, color=SLATE)
        text_box(s, x=x + 0.30, y=y + cell_h - 0.28, w=cell_w - 0.5, h=0.28,
                 text=src2, size=10, italic=True, color=SLATE)
    # Counter-fact gold strip
    gold_callout(s, 0.55, 6.05, 12.25, 0.90,
                 "Контр-факт: ~90% AI-пилотов в РФ не доходят до прода. CNews / Vedomosti / Intellectual Analytics, март 2026. [5]",
                 size=14)
    refs_of_slide(s, "s08")
    notes_with_sources(s, "s08")


def build_s09(p):
    """4 breakthroughs 2023-2026 — issue #153 fix #7 (episode 4 replaced).

    Episode 4: Kimi K2.5 (Moonshot) → Georgi Gerganov / llama.cpp / ggml.ai.
    Solo project → joined Hugging Face Feb 2026 (kept full autonomy) →
    100K+ GitHub stars March 2026, faster than PyTorch/TensorFlow.
    Lesson differentiation vs episode 3 (OpenClaw): ep.3 = one person can
    ship a product/agent that moves markets in weeks (top-down product);
    ep.4 = one person can ship infrastructure that becomes the backbone
    of the whole open ecosystem (bottom-up enabling layer).
    """
    s = blank(p)
    slide_title(s, "Пространство открыто: 4 прорыва 2023–2026 от не-первых игроков.", size=26)
    episodes = [
        ("сентябрь\n2023", "Mistral 7B",
            "Apache 2.0\nобходит Llama-2 13B",
            "Mistral AI (FR) [1]", MID, False),
        ("январь\n2025", "DeepSeek R1",
            "$589B\nNvidia drop за день",
            "DeepSeek (CN) [2, 3, 4]", GOLD, True),
        ("ноябрь\n2025", "OpenClaw",
            "100K★ stars\nза квартал",
            "P. Steinberger [5]", MID, False),
        ("февраль\n2026", "llama.cpp",
            "100K+★ на GitHub\nбыстрее PyTorch",
            "G. Gerganov / ggml.ai [6]", LIGHT, False),
    ]
    card_y = 2.05
    card_w = 2.95
    card_h = 4.0
    gap = 0.25
    start_x = (SLIDE_W_IN - (card_w * 4 + gap * 3)) / 2.0
    for i, (date, name, fact, org, color, is_gold) in enumerate(episodes):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h,
                  stroke=color, stroke_pt=2.5 if is_gold else 1.5)
        # Date band top
        filled_rect(s, x, card_y, card_w, 0.7, color, radius=True, radius_adj=0.1)
        text_box(s, x=x, y=card_y + 0.06, w=card_w, h=0.6, text=date,
                 size=14, bold=True, color=DEEP if color == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        # Name big
        text_box(s, x=x + 0.15, y=card_y + 0.95, w=card_w - 0.3, h=0.7, text=name,
                 size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.10)
        # Main fact
        text_box(s, x=x + 0.20, y=card_y + 1.85, w=card_w - 0.4, h=1.4, text=fact,
                 size=14, bold=is_gold, color=color if is_gold else DEEP,
                 align=PP_ALIGN.CENTER, line_spacing=1.30)
        # Org
        text_box(s, x=x, y=card_y + card_h - 0.55, w=card_w, h=0.4, text=org,
                 size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.42, 12.25, 0.55,
                 "Не отчаивайтесь: серьёзные прорывы делают разные команды. Курс — про устойчивые концепты, переживающие смену поколений моделей.",
                 size=13)
    refs_of_slide(s, "s09")
    notes_with_sources(s, "s09")


def build_s10(p):
    """Section 3 divider — zoom-in state of unified nav template (Fix-17).

    Same 6-card grid as s02a/s22/s27; card 3 is gold-FILLED with white text.
    Title encodes which section we're entering; frame phrase is the section's
    one-sentence framing.
    """
    s = blank(p)
    # Fix-19: sub_marker removed — gold-filled active card is sole navigation indicator.
    nav_slide(s, here_idx=3,
              title="Раздел 3 · Четыре способа реализации систем с AI",
              frame_phrase="Не альтернативы, а слои.")
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """Layers not alternatives — Fix-9 v2 (Phase 12.6, 2026-05-13).

    Layers stack BOTTOM-ALIGNED with shared bottom edge (each outer layer
    extends only upward) — more vertical room at top of each layer for
    its component caption strip on a teal-tinted background.
    Outer box must clear the slide title (top edge ≥ 1.65").
    """
    s = blank(p)
    slide_title(s, "Способы реализации систем с AI: не альтернативы, а слои.", size=26)
    # Bottom-aligned nested layers — common bottom edge at 6.65 (above 7.05 takeaway).
    # Outer "Приложение" must have top edge ≥ 1.65 (clear of title) — so outer h ≤ 5.0.
    by_base = 6.65
    cx = 9.20  # right-half centre
    sizes = [
        # (w, h, color, label, components-string)
        (7.6, 5.0, DEEP,  "Приложение",
            "+ AI внутри продукта · формы, кнопки, интеграции"),
        (5.9, 3.6, MID,   "Агент",
            "+ инструменты (API, поиск, код) · планирование · vector DB"),
        (4.3, 2.3, LIGHT, "Чат",
            "+ UI диалога · память истории сообщений"),
        (2.6, 1.0, TEAL,  "Модель",
            "stateless: вход → модель → выход"),
    ]
    # Draw rings — outer first
    for (w, h, color, _label, _comp) in sizes:
        x = cx - w / 2
        y = by_base - h
        ocean_box(s, x, y, w, h, fill=WHITE, stroke=color, stroke_pt=2.5)
    # Per-layer FILLED label strip in the TOP of each ring (Fix-9: better visibility)
    for i, (w, h, color, label, comp) in enumerate(sizes):
        x = cx - w / 2
        y = by_base - h
        if i < 3:
            strip_h = 0.62
            strip_pad = 0.16
            filled_rect(s, x + strip_pad, y + strip_pad - 0.04, w - 2 * strip_pad, strip_h,
                        TEAL_TINT, stroke=color, stroke_pt=1.0, radius=True, radius_adj=0.20)
            text_box(s, x=x + strip_pad + 0.18, y=y + strip_pad - 0.02, w=w - 2 * strip_pad - 0.36, h=0.26,
                     text=label, size=13, bold=True, color=color, align=PP_ALIGN.LEFT)
            text_box(s, x=x + strip_pad + 0.18, y=y + strip_pad + 0.24, w=w - 2 * strip_pad - 0.36, h=0.30,
                     text=comp, size=11, italic=True, color=DEEP,
                     align=PP_ALIGN.LEFT, line_spacing=1.18)
        else:
            filled_rect(s, x + 0.10, y + 0.10, w - 0.20, h - 0.20, TEAL_TINT,
                        stroke=color, stroke_pt=1.0, radius=True, radius_adj=0.20)
            text_box(s, x=x, y=y + 0.13, w=w, h=0.30, text=label,
                     size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
            text_box(s, x=x, y=y + 0.46, w=w, h=0.30, text=comp,
                     size=10, italic=True, color=DEEP, align=PP_ALIGN.CENTER)
    # Left explanation column (restored — was missing in iter9)
    text_box(s, x=0.55, y=2.0, w=4.6, h=0.55,
             text="Каждый следующий слой",
             size=18, bold=True, color=DEEP)
    text_box(s, x=0.55, y=2.55, w=4.6, h=0.55,
             text="включает предыдущий",
             size=18, bold=True, color=DEEP)
    text_box(s, x=0.55, y=3.40, w=4.6, h=2.6,
             text="Не четыре альтернативные технологии, а четыре способа реализации одной задачи. [1] Каждое следующее обёртывание добавляет способности — и сложность, стоимость, потенциал ошибок. [2]",
             size=13, color=DEEP, line_spacing=1.40)
    gold_callout(s, 0.55, 6.10, 4.6, 0.85,
                 "Выбор слоя — инженерное решение, не альтернатива.",
                 size=12)
    refs_of_slide(s, "s11")
    notes_with_sources(s, "s11")


def build_s12(p):
    """Classification matrix — issue #153 fix #8: readability pass.

    - Lucide icons per task (tag/scan/search/sparkles/trending-up/list-checks).
    - Single-line column headers (no «Класси-/фикация» wraps).
    - Matrix filled with concrete examples in most cells (≥3/4 coverage).
    - Axis headers enlarged (11pt → 14pt bold) for readability per issue #153.
    - YOLO gold-highlight NEUTRALIZED to a regular cell — focus is axis
      readability, not a callback accent competing for attention.
    """
    s = blank(p)
    slide_title(s, "Классификация AI-систем [1] — две оси: тип задачи × модальность. [2]", size=26)
    matrix_x, matrix_y = 0.55, 1.65
    matrix_w, matrix_h = 12.25, 5.05
    ocean_box(s, matrix_x, matrix_y, matrix_w, matrix_h)
    # Task (X axis) headers — short single-line labels, enlarged (fix #8)
    tasks = [
        ("Классиф.",  "lucide-tag-blue.png"),
        ("Распозн.",  "lucide-scan-line-blue.png"),
        ("Поиск",     "lucide-search-blue.png"),
        ("Генерац.",  "lucide-sparkles-blue.png"),
        ("Прогноз",   "lucide-trending-up-blue.png"),
    ]
    modalities = ["Текст", "Изображ.", "Звук / видео", "Структ. данные"]
    grid_left = matrix_x + 1.55
    grid_top = matrix_y + 1.10
    grid_w = matrix_w - 1.75
    grid_h = matrix_h - 1.35
    cell_w = grid_w / len(tasks)
    cell_h = grid_h / len(modalities)
    # Column headers — icon above label, enlarged bold for readability (fix #8)
    for i, (label, icon) in enumerate(tasks):
        x = grid_left + i * cell_w
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon,
                      x=x + cell_w / 2 - 0.20, y=matrix_y + 0.18, w=0.40, h=0.40)
        text_box(s, x=x, y=matrix_y + 0.60, w=cell_w, h=0.40, text=label,
                 size=14, bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.10)
    # Row headers — enlarged (fix #8)
    for j, m in enumerate(modalities):
        y = grid_top + j * cell_h
        text_box(s, x=matrix_x + 0.10, y=y + cell_h / 2 - 0.17, w=1.40, h=0.44,
                 text=m, size=14, bold=True, color=MID, align=PP_ALIGN.RIGHT)
    # Grid lines
    for i in range(len(tasks) + 1):
        x = grid_left + i * cell_w
        filled_rect(s, x - 0.005, grid_top, 0.01, grid_h, SOFT_GREY)
    for j in range(len(modalities) + 1):
        y = grid_top + j * cell_h
        filled_rect(s, grid_left, y - 0.005, grid_w, 0.01, SOFT_GREY)
    # Cells filled with examples — (task_col, modality_row, label, color)
    # Color map: TEAL = generation, MID = recognition/search/classification,
    # LIGHT = forecast/planning. Fix #8: YOLO no longer gold-highlighted —
    # regular cell like its neighbors, axis readability is the focus here.
    cells = [
        # Классификация column
        (0, 0, "BERT, спам", MID),
        (0, 1, "ResNet", MID),
        (0, 2, "PANNs", MID),
        (0, 3, "XGBoost", MID),
        # Распознавание column
        (1, 0, "spaCy NER", LIGHT),
        (1, 1, "YOLO", LIGHT),
        (1, 2, "Whisper", LIGHT),
        (1, 3, "OCR таблиц", LIGHT),
        # Поиск column
        (2, 0, "BM25", MID),
        (2, 1, "CLIP", MID),
        (2, 2, "Shazam", MID),
        (2, 3, "vector DB", MID),
        # Генерация column
        (3, 0, "GPT-4o, Claude", TEAL),
        (3, 1, "DALL-E, MJ", TEAL),
        (3, 2, "ElevenLabs, Sora", TEAL),
        (3, 3, "Codex tab.", TEAL),
        # Прогноз column — #179: text×forecast now shares the "next-token
        # prediction" framing with generation (same models, same color).
        (4, 0, "GPT-4o, Claude", TEAL),
        (4, 1, "прогноз кадра", LIGHT),
        (4, 2, "прогноз видео", LIGHT),
        (4, 3, "Prophet, ARIMA", LIGHT),
    ]
    for ti, mi, label, color in cells:
        x = grid_left + ti * cell_w + 0.06
        y = grid_top + mi * cell_h + 0.08
        cw = cell_w - 0.12
        ch = cell_h - 0.16
        if label == "—":
            text_box(s, x=x, y=y + ch/2 - 0.15, w=cw, h=0.30, text=label,
                     size=11, color=SOFT_GREY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            continue
        filled_rect(s, x, y, cw, ch, color, radius=True, radius_adj=0.18)
        text_box(s, x=x + 0.04, y=y + ch/2 - 0.16, w=cw - 0.08, h=0.34, text=label,
                 size=12, bold=False, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
    # Bottom note — Gold reserved for the ≥1×/slide highlight rule (YOLO example, not cell color)
    gold_callout(s, 0.55, 6.47, 12.25, 0.50,
                 "YOLO — детекция объектов из демо в начале лекции. Подход к обучению и архитектура — позже в курсе.",
                 size=12)
    refs_of_slide(s, "s12")
    notes_with_sources(s, "s12")


def build_s13(p):
    """Control quadrant 2x2 — model/chat/agent placement.

    Fix-16 (2026-05-13): Axes swapped so Agent = right-top quadrant.
    X axis = «Делегирование от пользователя» (низкий → высокий)
    Y axis = «Контроль разработчика» (низкий → высокий)
    Diagonal: Модель (left-bottom) → Чат (center) → Агент (right-top, gold).
    Empty corners labelled: «нет смысла» (left-top), «опасная зона» (right-bottom).
    """
    s = blank(p)
    slide_title(s, "Одна задача, три способа: контроль распределяется между разработчиком и пользователем.", size=22)
    # Quadrant area — shrunk vertically a bit to leave room for axis label + callout below
    qx, qy = 1.95, 1.9
    qw, qh = 7.15, 3.95
    # Box outline
    filled_rect(s, qx, qy, qw, qh, WHITE, stroke=LIGHT, stroke_pt=1.5, radius=True, radius_adj=0.04)
    # Internal cross lines
    filled_rect(s, qx, qy + qh / 2 - 0.005, qw, 0.01, SOFT_GREY)
    filled_rect(s, qx + qw / 2 - 0.005, qy, 0.01, qh, SOFT_GREY)
    # Y axis label (left, vertical conceptual) — Fix-16: «Контроль разработчика»
    text_box(s, x=qx - 1.9, y=qy + qh / 2 - 0.40, w=1.8, h=0.85,
             text="Контроль\nразработчика",
             size=14, bold=True, color=MID, align=PP_ALIGN.RIGHT, line_spacing=1.18)
    # ↑ arrow + «высокий» at top of Y axis (just outside quadrant, near top-left corner)
    text_box(s, x=qx - 1.55, y=qy - 0.10, w=1.45, h=0.32, text="высокий ↑",
             size=14, bold=True, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # «низкий» — at bottom-left of Y axis area, OUTSIDE quadrant
    text_box(s, x=qx - 1.55, y=qy + qh + 0.06, w=1.45, h=0.32, text="низкий ↓",
             size=14, bold=True, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # X axis label (bottom) — Fix-16: «Делегирование от пользователя»
    text_box(s, x=qx + qw / 2 - 2.4, y=qy + qh + 0.08, w=4.8, h=0.40,
             text="Делегирование от пользователя",
             size=15, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # Arrow + range markers — Fix-16: moved BELOW X-axis label (outside quadrant) so they don't collide with
    # Agent at top-right or Model at bottom-left circles inside the quadrant.
    text_box(s, x=qx + 0.05, y=qy + qh + 0.48, w=1.5, h=0.30, text="← низкий",
             size=13, bold=True, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    text_box(s, x=qx + qw - 1.55, y=qy + qh + 0.48, w=1.5, h=0.30, text="высокий →",
             size=13, bold=True, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # Empty-quadrant labels — italic small grey, near corners (Fix-16)
    # Top-left (X=low delegation, Y=high разраб control) — «нет смысла»
    text_box(s, x=qx + 0.15, y=qy + 0.10, w=2.0, h=0.32,
             text="нет смысла", size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # Bottom-right (X=high delegation, Y=low разраб control) — «опасная зона»
    text_box(s, x=qx + qw - 2.15, y=qy + qh - 0.42, w=2.0, h=0.32,
             text="опасная зона", size=10, italic=True, color=SLATE, align=PP_ALIGN.RIGHT)
    # Three points on quadrant — Fix-16 placement.
    # PPTX coords: small fy = top of slide; large fy = bottom.
    # Y axis "high" is visually at top → high control разраб = small fy.
    # Strategy: corner circles + sub-text BELOW fit cleanly inside their respective half.
    # Quadrant: qy=1.9 to qy+qh=5.85. Cross-line at qy + qh/2 = 3.875.
    # Circle d=0.95 (slightly smaller than v3.0's 1.05 to free space for sub-text).
    # issue #153 fix #9: per-point sub-labels refined — short (≤6 words),
    # concrete «что характерно для решения этим способом» applied to the
    # PDF-contract task from speaker notes.
    pts = [
        # (fx, fy, label, sub, color, is_gold)
        # Модель: bottom-left. fy=0.68 → cy=4.586, circle 4.11-5.06, sub 5.18-5.78 (in bottom half).
        (0.20, 0.68, "Модель", "Сам интегрирует API, полный контроль", LIGHT, False),
        # Чат: center.
        (0.50, 0.50, "Чат", "Диалог, уточнения по ходу", MID, False),
        # Агент: top-right. fy=0.20 → cy=2.69, circle 2.21-3.16, sub 3.29-3.74 (in top half).
        (0.80, 0.20, "Агент", "Делегирование целиком, решает оркестратор", GOLD, True),
    ]
    for fx, fy, label, sub, color, is_gold in pts:
        cx = qx + fx * qw
        cy = qy + fy * qh
        d = 0.95
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - d/2), Inches(cy - d/2),
                                 Inches(d), Inches(d))
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.color.rgb = DEEP; shp.line.width = Pt(2.0 if is_gold else 1.0)
        disable_shadow(shp)
        text_box(s, x=cx - 0.6, y=cy - 0.20, w=1.2, h=0.4, text=label,
                 size=14, bold=True, color=DEEP if is_gold else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Sub-text always BELOW circle. issue #155 fix #183/#184: labels are
        # now single-line (no \n). iter-2: 2.6"/9pt still wrapped for the two
        # longest labels (word_wrap=True in text_box) — widened to 3.4" and
        # reduced to 8.5pt. Agent's box (cx=7.67) stays clear of the task-box
        # at qx+qw+0.5=9.6 (right edge 7.67+1.7=9.37); Модель's box (cx=3.38)
        # stays clear of the left slide margin (left edge 3.38-1.7=1.68).
        text_box(s, x=cx - 1.70, y=cy + 0.50, w=3.4, h=0.35, text=sub,
                 size=8.5, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    # Right side — fixed task box
    task_x = qx + qw + 0.5
    task_w = SLIDE_W_IN - task_x - 0.35
    ocean_box(s, task_x, qy + 0.5, task_w, 3.3)
    text_box(s, x=task_x + 0.25, y=qy + 0.65, w=task_w - 0.5, h=0.4,
             text="Одна и та же задача", size=14, bold=True, color=TEAL)
    text_box(s, x=task_x + 0.25, y=qy + 1.10, w=task_w - 0.5, h=2.0,
             text="Извлечь поля из входящего PDF-договора:\n• дата подписания\n• контрагент\n• сумма\n• срок действия\n\nи положить в таблицу.",
             size=13, color=DEEP, line_spacing=1.40)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.42, 12.25, 0.55,
                 "Распределение контроля — инженерное решение, а не достоинство одного способа над другим. [1]",
                 size=13)
    refs_of_slide(s, "s13")
    notes_with_sources(s, "s13")


# build_s14 (mini-divider «Разберём подробнее») deleted under Fix-17, 2026-05-13.
# Reason: paraphrased s10 framing; 4-type icons broke the lecture's 5-section
# navigation grammar. Verbal transition moved to s13 «Лектору» notes.


def build_s15(p):
    """Model with pipeline schema — issue #153 fix #11.

    - Eyebrow pill «МОДЕЛЬ» (fix #10, consistent across s15-s19a).
    - Alignment pass: owner labels now centred UNDER each block (were
      offset from a stale start_x reference) and pipeline vertically
      re-centred to clear the eyebrow pill.
    - Explicit outer frame around the whole 5-block pipeline labelled
      «Это уже приложение» — key idea: model = one component, the WHOLE
      wired-up pipeline is already an application.
    """
    s = blank(p)
    eyebrow_pill(s, "МОДЕЛЬ")
    slide_title(s, "Модель — компонент, не система. Инференс: вход → препроцессинг → модель → постпроцессинг → выход.", size=20, y=0.85)
    # Outer framing box — the whole pipeline IS an application (key idea, fix #11)
    frame_x, frame_y = 0.75, 2.05
    frame_w, frame_h = 11.85, 2.55
    filled_rect(s, frame_x, frame_y, frame_w, frame_h, WHITE, stroke=GOLD, stroke_pt=2.0,
                radius=True, radius_adj=0.05)
    text_box(s, x=frame_x + 0.15, y=frame_y - 0.32, w=6.0, h=0.35,
             text="Это уже приложение", size=13, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    # Horizontal pipeline — centred INSIDE the outer frame (fix #11 alignment)
    pip_y = frame_y + 0.55
    pip_h = 1.35
    blocks = [
        ("Сырой\nвход", "кадр камеры,\nтекст, звук", LIGHT),
        ("Препро-\nцессинг", "масштабирование,\nобрезка, токенизация", LIGHT),
        ("Модель", "инференс", MID),
        ("Постпро-\nцессинг", "фильтрация,\nнормализация", LIGHT),
        ("Выход", "JSON, метка,\nдействие", LIGHT),
    ]
    n = len(blocks)
    block_w = 1.95
    arrow_w = 0.50
    total_w = block_w * n + arrow_w * (n - 1)
    start_x = (SLIDE_W_IN - total_w) / 2.0
    for i, (name, sub, color) in enumerate(blocks):
        x = start_x + i * (block_w + arrow_w)
        is_model = (i == 2)
        filled_rect(s, x, pip_y, block_w, pip_h, color,
                    stroke=DEEP if is_model else None, stroke_pt=2.0 if is_model else 0.0,
                    radius=True, radius_adj=0.15)
        text_box(s, x=x, y=pip_y + 0.14, w=block_w, h=0.50, text=name,
                 size=14 if not is_model else 16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.10)
        text_box(s, x=x + 0.08, y=pip_y + 0.72, w=block_w - 0.16, h=0.58, text=sub,
                 size=9.5, italic=True, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.20)
        if i < n - 1:
            # Fix-11: Use proper RIGHT_ARROW shape (not segmented rect+triangle).
            ax = x + block_w + 0.04
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(ax), Inches(pip_y + pip_h / 2 - 0.20),
                                       Inches(arrow_w - 0.08), Inches(0.40))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
            disable_shadow(arrow)
    # Owner labels — aligned exactly under each block (fix #11: fixed offset bug)
    owner_labels = ["внешняя система", "разработчик", "AI-модель", "разработчик", "приложение"]
    for i, label in enumerate(owner_labels):
        x = start_x + i * (block_w + arrow_w)
        is_model = (i == 2)
        text_box(s, x=x, y=pip_y + pip_h + 0.08, w=block_w, h=0.32,
                 text=f"↑ {label}", size=10, italic=True,
                 color=DEEP if is_model else (MID if i in (1, 3) else SLATE),
                 align=PP_ALIGN.CENTER, bold=(i != 0 and i != 4))
    # 4 model examples
    examples = [
        ("YOLOv8", "детекция на изображениях"),
        ("Whisper", "распознавание речи"),
        ("Stable Diffusion", "генерация изображений"),
        ("AlphaFold", "прогноз структур белков [2]"),
    ]
    ex_y = 5.05
    ex_w = 2.8
    ex_h = 1.25
    ex_gap = 0.20
    ex_start_x = (SLIDE_W_IN - (ex_w * 4 + ex_gap * 3)) / 2.0
    for i, (name, role) in enumerate(examples):
        x = ex_start_x + i * (ex_w + ex_gap)
        ocean_box(s, x, ex_y, ex_w, ex_h)
        text_box(s, x=x, y=ex_y + 0.18, w=ex_w, h=0.42, text=name,
                 size=15, bold=True, color=MID, align=PP_ALIGN.CENTER)
        text_box(s, x=x, y=ex_y + 0.62, w=ex_w, h=0.58, text=role,
                 size=11, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.25)
    gold_callout(s, 0.55, 6.42, 12.25, 0.55,
                 "Препроцессинг и постпроцессинг — ответственность разработчика. [1] YOLO = 50 строк; рабочая система с YOLO = сотни строк.",
                 size=12)
    refs_of_slide(s, "s15")
    notes_with_sources(s, "s15")


def build_s16(p):
    """Chat cycle — Fix-12 (Phase 12.6, 2026-05-13): compact dialog-cycle visual.

    Replaces the 6-step linear flow with a more intuitive dialog-cycle layout:
       [system prompt]
              ↓ merges
       [User] → [сообщение] ──→ [LLM]
                                  │
       [User] ← [ответ]  ←────────┘
                  ↓
                  ⋮  next iteration

    + 2 gold callouts on the right (system prompt control / context window).
    + Bottom takeaway preserved.
    """
    s = blank(p)
    eyebrow_pill(s, "ЧАТ")
    slide_title(s, "Как работает чат: цикл диалога.", size=28, y=0.85)

    # ─── Layout constants ───
    # Left column for the dialog-cycle visual; right column for callouts.
    diag_x0 = 0.55
    diag_w = 7.30          # dialog area width
    diag_y0 = 2.05
    diag_h = 4.65          # dialog area height

    # User icon column (x), message column (x), LLM box column (x)
    user_x = diag_x0 + 0.05
    user_w = 0.95
    msg_x = diag_x0 + 1.15
    msg_w = 3.75
    llm_x = diag_x0 + 5.05
    llm_w = 2.20

    # Vertical positions
    sysprompt_y = diag_y0 + 0.05      # system prompt sits above [сообщение]
    sysprompt_h = 0.65
    msg_y = diag_y0 + 0.95            # «сообщение» row
    row_h = 0.85
    answer_y = diag_y0 + 2.30         # «ответ» row, BELOW «сообщение»
    dots_y = diag_y0 + 3.40           # «⋮» continuation indicator

    # ─── User icon (twice — once per row, sharing a small label) ───
    for label_y in (msg_y, answer_y):
        # Round filled circle for user
        ud = 0.85
        uy = label_y + (row_h - ud) / 2
        ucirc = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(user_x), Inches(uy),
                                   Inches(ud), Inches(ud))
        ucirc.fill.solid(); ucirc.fill.fore_color.rgb = MID
        ucirc.line.fill.background()
        disable_shadow(ucirc)
        text_box(s, x=user_x, y=uy + 0.18, w=ud, h=ud - 0.30,
                 text="USER", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ─── System prompt block, sits ABOVE the message row ───
    filled_rect(s, msg_x, sysprompt_y, msg_w, sysprompt_h, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.2, radius=True, radius_adj=0.22)
    text_box(s, x=msg_x + 0.20, y=sysprompt_y + 0.08, w=msg_w - 0.40, h=0.30,
             text="Системный промпт",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=msg_x + 0.20, y=sysprompt_y + 0.36, w=msg_w - 0.40, h=0.30,
             text="роль, ограничения, формат ответа",
             size=10, italic=True, color=DEEP, align=PP_ALIGN.LEFT)

    # ─── Message box (USER → LLM) ───
    filled_rect(s, msg_x, msg_y, msg_w, row_h, WHITE, stroke=MID, stroke_pt=1.5,
                radius=True, radius_adj=0.18)
    text_box(s, x=msg_x + 0.20, y=msg_y + 0.10, w=msg_w - 0.40, h=0.30,
             text="Сообщение",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=msg_x + 0.20, y=msg_y + 0.42, w=msg_w - 0.40, h=0.36,
             text="вопрос / инструкция / фрагмент данных",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # Merge-arrow from system prompt down into message (visual «объединяются»)
    merge_x = msg_x + msg_w / 2 - 0.18
    merge_y = sysprompt_y + sysprompt_h + 0.02
    merge_arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(merge_x), Inches(merge_y),
                                     Inches(0.36), Inches(0.24))
    merge_arrow.fill.solid(); merge_arrow.fill.fore_color.rgb = GOLD
    merge_arrow.line.fill.background()
    disable_shadow(merge_arrow)

    # ─── USER → message arrow (right) ───
    arr_y = msg_y + row_h / 2 - 0.10
    fwd = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             Inches(user_x + user_w + 0.02), Inches(arr_y),
                             Inches(msg_x - (user_x + user_w + 0.02) - 0.02), Inches(0.20))
    fwd.fill.solid(); fwd.fill.fore_color.rgb = MID
    fwd.line.fill.background()
    disable_shadow(fwd)

    # ─── message → LLM arrow ───
    msg_to_llm = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(msg_x + msg_w + 0.02), Inches(arr_y),
                                    Inches(llm_x - (msg_x + msg_w + 0.02) - 0.02), Inches(0.20))
    msg_to_llm.fill.solid(); msg_to_llm.fill.fore_color.rgb = MID
    msg_to_llm.line.fill.background()
    disable_shadow(msg_to_llm)

    # ─── LLM box (spans both rows vertically) ───
    llm_h = (answer_y + row_h) - msg_y
    filled_rect(s, llm_x, msg_y, llm_w, llm_h, MID, radius=True, radius_adj=0.18,
                stroke=DEEP, stroke_pt=2.0)
    text_box(s, x=llm_x, y=msg_y + llm_h / 2 - 0.40, w=llm_w, h=0.45,
             text="LLM",
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=llm_x, y=msg_y + llm_h / 2 + 0.05, w=llm_w, h=0.40,
             text="модель",
             size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ─── LLM → answer arrow (left, going BACK toward USER row) ───
    # Vertical drop from LLM bottom-mid to answer row right edge
    llm_bot_y = msg_y + llm_h - 0.05
    # Diagonal-ish: short vertical at LLM bottom, then horizontal LEFT to answer box
    # We use a single LEFT_ARROW from LLM region to answer's right edge.
    ans_y_arr = answer_y + row_h / 2 - 0.10
    back = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                              Inches(msg_x + msg_w + 0.02), Inches(ans_y_arr),
                              Inches(llm_x - (msg_x + msg_w + 0.02) - 0.02), Inches(0.20))
    back.fill.solid(); back.fill.fore_color.rgb = MID
    back.line.fill.background()
    disable_shadow(back)

    # ─── Answer box ───
    filled_rect(s, msg_x, answer_y, msg_w, row_h, WHITE, stroke=MID, stroke_pt=1.5,
                radius=True, radius_adj=0.18)
    text_box(s, x=msg_x + 0.20, y=answer_y + 0.10, w=msg_w - 0.40, h=0.30,
             text="Ответ",
             size=13, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    text_box(s, x=msg_x + 0.20, y=answer_y + 0.42, w=msg_w - 0.40, h=0.36,
             text="токен за токеном · дописывается в историю",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)

    # ─── Answer → USER arrow (left) ───
    ans_to_user = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                     Inches(user_x + user_w + 0.02), Inches(ans_y_arr),
                                     Inches(msg_x - (user_x + user_w + 0.02) - 0.02), Inches(0.20))
    ans_to_user.fill.solid(); ans_to_user.fill.fore_color.rgb = MID
    ans_to_user.line.fill.background()
    disable_shadow(ans_to_user)

    # ─── Accumulating-history visual (issue #153 fix #12) ───
    # Replaces the old «⋮ следующая итерация» hint with an explicit growing
    # block showing the model re-reads the WHOLE history every single step —
    # not an increment. 4 stacked segments, growing width left→right.
    hist_label_y = dots_y - 0.05
    text_box(s, x=msg_x, y=hist_label_y, w=msg_w, h=0.28,
             text="Следующий шаг — весь текст заново:",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    hist_bar_y = hist_label_y + 0.30
    hist_bar_h = 0.24
    seg_labels = ["сист. промпт", "сообщение 1", "ответ 1", "сообщение 2"]
    seg_colors = [GOLD, MID, TEAL, MID]
    seg_w_unit = msg_w / len(seg_labels)
    cur_x = msg_x
    for i, (seg_label, seg_color) in enumerate(zip(seg_labels, seg_colors)):
        seg_w = seg_w_unit * (i + 1) / len(seg_labels) + seg_w_unit * 0.5
        seg_w = min(seg_w, msg_w - (cur_x - msg_x))
        filled_rect(s, cur_x, hist_bar_y, seg_w_unit - 0.03, hist_bar_h, seg_color,
                    radius=True, radius_adj=0.3)
        cur_x += seg_w_unit
    text_box(s, x=msg_x, y=hist_bar_y + hist_bar_h + 0.06, w=msg_w, h=0.30,
             text="весь текст заново на каждом шаге — не инкремент",
             size=10.5, italic=True, bold=True, color=DEEP, align=PP_ALIGN.LEFT)

    # ─── Right column: 2 callouts ───
    cb_x = diag_x0 + diag_w + 0.30
    cb_w = SLIDE_W_IN - cb_x - 0.55
    gold_callout(s, cb_x, sysprompt_y, cb_w, 1.65,
                 "Контроль через системный промпт.\n"
                 "Промпт задаёт роль, ограничения, формат. Один и тот же "
                 "чат настраивается под разные сценарии — это инженерный "
                 "рычаг разработчика. [1]",
                 size=12)
    gold_callout(s, cb_x, sysprompt_y + 1.95, cb_w, 1.85,
                 "Ограничение — контекстное окно.\n"
                 "128k–1M токенов; когда история перестаёт помещаться, "
                 "старые сообщения выпадают и чат «забывает» начало "
                 "длинного разговора. Это ограничение модели, не баг.",
                 size=12)

    # ─── Bottom takeaway — issue #153 fix #12: «а не магия» tail removed ───
    text_box(s, x=0.55, y=6.62, w=12.25, h=0.35,
             text="Чат — это конвейер «собрать → подать → дописать → показать».",
             size=13, italic=True, bold=True, color=DEEP, align=PP_ALIGN.CENTER)

    refs_of_slide(s, "s16")
    notes_with_sources(s, "s16")


def build_s17(p):
    """Chat = model + UI + memory; case + LLM bar chart."""
    s = blank(p)
    eyebrow_pill(s, "ЧАТ")
    slide_title(s, "Чат = модель + интерфейс + память диалога.", size=28, y=0.85)
    # Left: case card
    case_x, case_y, case_w, case_h = 0.55, 2.15, 6.5, 4.5
    ocean_box(s, case_x, case_y, case_w, case_h)
    text_box(s, x=case_x + 0.30, y=case_y + 0.25, w=case_w - 0.6, h=0.4,
             text="Кейс — типовой для чата", size=14, bold=True, color=TEAL)
    text_box(s, x=case_x + 0.30, y=case_y + 0.7, w=case_w - 0.6, h=0.95,
             text="Инженер получил непонятное ТЗ от смежного отдела, нужно составить чек-лист своей работы.",
             size=15, color=DEEP, line_spacing=1.30)
    # Mock dialog
    dlg_y = case_y + 1.85
    dlg_h = 2.0
    filled_rect(s, case_x + 0.4, dlg_y, case_w - 0.8, dlg_h, WHITE,
                stroke=SOFT_GREY, stroke_pt=1.0, radius=True, radius_adj=0.08)
    text_runs(s, case_x + 0.6, dlg_y + 0.15, case_w - 1.2, dlg_h - 0.3, [
        {"text": "Вы:  ", "size": 12, "bold": True, "color": MID},
        {"text": "Объясни пункт 4.2 ТЗ простыми словами.", "size": 12, "color": DEEP},
        {"newpara": True, "text": "Чат:  ", "size": 12, "bold": True, "color": TEAL},
        {"text": "Это требование к…", "size": 12, "color": DEEP},
        {"newpara": True, "text": "Вы:  ", "size": 12, "bold": True, "color": MID},
        {"text": "Составь чек-лист на 5 пунктов.", "size": 12, "color": DEEP},
        {"newpara": True, "text": "Чат:  ", "size": 12, "bold": True, "color": TEAL},
        {"text": "1. Проверить…  2. Согласовать…  3. …", "size": 12, "color": DEEP},
    ], line_spacing=1.45)
    text_box(s, x=case_x + 0.30, y=case_y + case_h - 0.5, w=case_w - 0.6, h=0.35,
             text="Разовая задача с уточнениями — оптимально для чата. Не модель, не агент, не приложение.",
             size=11, italic=True, color=LIGHT, line_spacing=1.30)
    # Right side — Fix-13 (Phase 12.6, 2026-05-13):
    # Replaced bar chart with «production-disclaimer» card explaining that
    # pure chats are rarely used in production — almost always wrapped with
    # an agent at minimum (memory, RAG, tool calls).
    disc_x, disc_y, disc_w, disc_h = case_x + case_w + 0.35, 2.15, 5.4, 4.5
    ocean_box(s, disc_x, disc_y, disc_w, disc_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    # issue #153 QA fix #2 (P2, Russification): "Disclaimer"/"production" →
    # Russian per §5.8 table ("production use" → "промышленное применение").
    text_box(s, x=disc_x + 0.30, y=disc_y + 0.25, w=disc_w - 0.6, h=0.4,
             text="Оговорка для промышленных систем", size=14, bold=True, color=DEEP)
    # issue #153 QA fix #2 follow-up: "промышленной эксплуатации" (25 chars)
    # is longer than "production" (10 chars) — the old fixed y+1.95 offset for
    # the paragraph below assumed a 2-line wrap; the longer RU text wraps to
    # 3 lines at 18pt and collided with the text below it. Fix: smaller font
    # (18→16pt) to encourage 2-line wrap, and push the next block down
    # (1.95→2.35) to clear the worst-case 3-line wrap height.
    text_box(s, x=disc_x + 0.30, y=disc_y + 0.85, w=disc_w - 0.6, h=1.4,
             text="Чистые чаты почти не используются в промышленной эксплуатации.",
             size=16, bold=True, color=DEEP, line_spacing=1.25)
    text_box(s, x=disc_x + 0.30, y=disc_y + 2.35, w=disc_w - 0.6, h=2.0,
             text="Почти везде они расширены до агентов [1] [2] — хотя бы для "
                  "долгосрочной памяти и поиска по корпоративной базе (RAG).\n\n"
                  "Архитектуру агента разберём на следующем слайде.",
             size=12, color=DEEP, line_spacing=1.40)
    # issue #153 fix #13: «Возвращаемся к...» removed — direct statement instead.
    gold_callout(s, 0.55, 6.47, 12.25, 0.50,
                 "Выбор чата — точка на шкале взаимодействия, не единственно верный вариант.",
                 size=13)
    refs_of_slide(s, "s17")
    notes_with_sources(s, "s17")


def build_s18(p):
    """Agent architecture — issue #153 fix #14: FULL REDESIGN.

    v3.2 hub-and-spoke schema (Chat center + Orchestrator above + Memory/Tools
    flanking) was flagged weak by the owner with no reference to fix toward.
    3 alternatives considered during the visual loop:
      (a) linear pipeline plan→act→observe→reflect with explicit loop-back —
          CHOSEN: clearest for a student with no prior architecture exposure,
          passes the 5-Second Test (main message = "it's a cycle of 4 named
          steps, and USER starts + receives the result").
      (b) cleaner hub-and-spoke (same layout, better drawn loop) — rejected,
          inherits the same "what's a hub-and-spoke architecture" cognitive
          load the v3.2 version already failed on.
      (c) sequence-diagram (USER→Orchestrator→Tool→LLM→USER swimlanes) —
          rejected, too much detail (swimlane crossings) for a 1.5-min slide;
          better suited for a deep-dive lecture, not this intro slide.
    Content unchanged: agent = chat + orchestrator + external memory + tools;
    ReAct cycle (Yao et al. 2022). Only the visual representation changes.
    Schema Readability Checklist (Architecture/Actor subtype): USER explicit
    (left, bidirectional labelled arrows) — PASS. Components grouped by tier
    (input / 4-stage loop / resources) — PASS. Connectors labelled — PASS.
    """
    s = blank(p)
    eyebrow_pill(s, "АГЕНТ")
    slide_title(s, "Агент = чат + оркестратор + внешняя память + инструменты. [1] [2] [3]", size=26, y=0.85)

    # ─── USER actor (left) ───
    # issue #155 QA fix #189: USER cluster was vertically offset below the
    # pipeline row (center 4.10" vs pipeline center 3.225") — raised so the
    # USER circle's vertical center matches the pipeline stage row center.
    user_d = 1.1
    user_x = 0.75
    stage_y_ref = 2.55
    stage_h_ref = 1.35
    user_y = stage_y_ref + stage_h_ref / 2 - user_d / 2
    ucirc = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(user_x), Inches(user_y),
                               Inches(user_d), Inches(user_d))
    ucirc.fill.solid(); ucirc.fill.fore_color.rgb = LIGHT
    ucirc.line.color.rgb = DEEP; ucirc.line.width = Pt(1.5)
    disable_shadow(ucirc)
    text_box(s, x=user_x, y=user_y + user_d/2 - 0.16, w=user_d, h=0.35,
             text="USER", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=user_x - 0.35, y=user_y + user_d + 0.06, w=user_d + 0.7, h=0.30,
             text="Пользователь", size=10.5, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

    # ─── 4-stage linear ReAct pipeline: План → Действие → Наблюдение → Рефлексия ───
    # issue #153 QA fix #1 (Russification, presentation-critic P1): stage labels
    # were bare English verbs (Plan/Act/Observe/Reflect) — translated to Russian
    # per §5.8, small English gloss kept in the sub-label for ReAct traceability.
    stages = [
        ("План", "план действий (Plan)", MID),
        ("Действие", "вызов инструмента (Act)", TEAL),
        ("Наблюдение", "результат в память (Observe)", LIGHT),
        ("Рефлексия", "цель достигнута? (Reflect)", MID),
    ]
    stage_y = 2.55
    stage_h = 1.35
    stage_w = 2.15
    arrow_w = 0.45
    n = len(stages)
    total_w = stage_w * n + arrow_w * (n - 1)
    start_x = user_x + user_d + 0.55
    for i, (name, sub, color) in enumerate(stages):
        x = start_x + i * (stage_w + arrow_w)
        filled_rect(s, x, stage_y, stage_w, stage_h, color, stroke=DEEP, stroke_pt=1.5,
                    radius=True, radius_adj=0.15)
        # issue #153 QA fix #1: Russian labels are longer than EN originals
        # («Наблюдение»/«Рефлексия» vs «Observe»/«Reflect») — smaller font for
        # long names keeps single-line fit inside the 2.15" stage box.
        name_size = 14 if len(name) > 8 else 17
        text_box(s, x=x, y=stage_y + 0.20, w=stage_w, h=0.45, text=name,
                 size=name_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text_box(s, x=x + 0.08, y=stage_y + 0.75, w=stage_w - 0.16, h=0.50, text=sub,
                 size=9.5, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.15)
        if i < n - 1:
            ax = x + stage_w + 0.03
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(ax), Inches(stage_y + stage_h/2 - 0.16),
                                       Inches(arrow_w - 0.06), Inches(0.32))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
            disable_shadow(arrow)
    end_x = start_x + total_w

    # USER → Plan (start) arrow
    a1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                            Inches(user_x + user_d + 0.05), Inches(stage_y + stage_h/2 - 0.10),
                            Inches(start_x - (user_x + user_d + 0.05) - 0.05), Inches(0.20))
    a1.fill.solid(); a1.fill.fore_color.rgb = DEEP
    a1.line.fill.background()
    disable_shadow(a1)

    # Reflect → USER (stop, result back) — L-shaped route: horizontal run
    # BELOW the pipeline boxes (unchanged height, so it doesn't cut across
    # the box labels), then a short vertical riser up to USER's own raised
    # center, landing an UP arrowhead right on USER's bottom edge.
    # issue #155 QA fix #191: original version's horizontal run stayed at the
    # old fixed height tied to the pipeline bottom, which (after #189 raised
    # USER higher) ended below the USER circle with a visible gap. A version
    # that instead ran the whole horizontal segment at USER's new center cut
    # straight through the pipeline box labels — also wrong. This L-route
    # keeps the horizontal leg below the boxes (clear of their text) and only
    # rises next to USER, where there is nothing else to collide with.
    stop_y = stage_y + stage_h + 0.55
    user_right_edge = user_x + user_d
    user_bottom = user_y + user_d
    filled_rect(s, end_x - 0.03, stage_y + stage_h, 0.06, stop_y - (stage_y + stage_h), TEAL)
    stop_arrow = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                    Inches(user_right_edge + 0.15), Inches(stop_y - 0.10),
                                    Inches(end_x - (user_right_edge + 0.15) - 0.03), Inches(0.20))
    stop_arrow.fill.solid(); stop_arrow.fill.fore_color.rgb = TEAL
    stop_arrow.line.fill.background()
    disable_shadow(stop_arrow)
    # Vertical riser from the horizontal run up to USER's bottom edge.
    riser_x = user_right_edge + 0.15
    riser_up = s.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                  Inches(riser_x - 0.10), Inches(user_bottom),
                                  Inches(0.20), Inches(stop_y - user_bottom))
    riser_up.fill.solid(); riser_up.fill.fore_color.rgb = TEAL
    riser_up.line.fill.background()
    disable_shadow(riser_up)
    # issue #153 QA fix #1 (Russification): "stop" → "стоп".
    # issue #155 QA fix #191 follow-up: label x-start nudged right of the new
    # riser arrow (riser occupies ~riser_x-0.10 .. riser_x+0.10) — previously
    # it started underneath the riser and had its first letter clipped.
    text_box(s, x=riser_x + 0.25, y=stop_y - 0.42, w=total_w - (riser_x + 0.25 - user_x), h=0.28,
             text="стоп → результат пользователю", size=10.5, italic=True, color=TEAL,
             align=PP_ALIGN.LEFT)

    # Reflect → Plan loop-back (continue) — gold arc above the pipeline
    loop_y = stage_y - 0.55
    loop_x0 = start_x + stage_w / 2
    loop_x1 = end_x - stage_w / 2
    filled_rect(s, loop_x0, loop_y, loop_x1 - loop_x0, 0.06, GOLD)
    filled_rect(s, loop_x0, loop_y, 0.06, stage_y - loop_y, GOLD)
    # issue #155 QA fix #190: flow is Reflect → (up into) loop-back bar → Plan,
    # so the vertical connector arrowhead must point UP (out of Рефлексия into
    # the bar), not DOWN (which read as the bar feeding INTO Рефлексия).
    left_arrow = s.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                    Inches(loop_x1 - 0.10), Inches(loop_y - 0.02),
                                    Inches(0.20), Inches(0.10 + (stage_y - loop_y)))
    left_arrow.fill.solid(); left_arrow.fill.fore_color.rgb = GOLD
    left_arrow.line.fill.background()
    disable_shadow(left_arrow)
    # issue #153 QA fix #1 (Russification): "continue" → "продолжить".
    text_box(s, x=loop_x0, y=loop_y - 0.35, w=loop_x1 - loop_x0, h=0.30,
             text="продолжить — цикл повторяется", size=11, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)

    # ─── Resources row below: Memory (under Observe) + Tools (under Act) ───
    res_y = stage_y + stage_h + 0.90
    res_h = 1.15
    tools_x = start_x + stage_w + arrow_w
    ocean_box(s, tools_x, res_y, stage_w, res_h, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=tools_x, y=res_y + 0.12, w=stage_w, h=0.35,
             text="ИНСТРУМЕНТЫ", size=11.5, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(s, x=tools_x + 0.10, y=res_y + 0.50, w=stage_w - 0.20, h=0.55,
             text="API, файлы,\ncode, search", size=10, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.20)
    mem_x = start_x + 2 * (stage_w + arrow_w)
    ocean_box(s, mem_x, res_y, stage_w, res_h, fill=SURFACE, stroke=LIGHT)
    text_box(s, x=mem_x, y=res_y + 0.12, w=stage_w, h=0.35,
             text="ПАМЯТЬ", size=11.5, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    text_box(s, x=mem_x + 0.10, y=res_y + 0.50, w=stage_w - 0.20, h=0.55,
             text="vector DB,\nфайлы, логи", size=10, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.20)
    # Connector lines: Act → Tools, Observe → Memory
    filled_rect(s, tools_x + stage_w/2 - 0.03, stage_y + stage_h, 0.06, res_y - (stage_y + stage_h), TEAL)
    filled_rect(s, mem_x + stage_w/2 - 0.03, stage_y + stage_h, 0.06, res_y - (stage_y + stage_h), LIGHT)
    # issue #153 QA fix #8 (P2, non-blocking): tiny "использует" caption on each
    # vertical connector for parity with the labelled horizontal flow above.
    # issue #155 QA fix P1-2 (3 sub-iterations):
    # (1) original mid-gap position (y = pipeline-bottom + half the gap to
    #     the resources row = 4.25") placed the label directly across BOTH
    #     the "стоп → результат пользователю" text row (4.03..4.31") AND the
    #     horizontal stop-arrow teal band (4.35..4.55"), cutting the "ь".
    # (2) moved label down to just below the stop-arrow band (stop_y+0.13 =
    #     4.58") — cleared the stop-arrow/text collision, but the label was
    #     still horizontally CENTERED on the vertical TEAL/LIGHT connector
    #     stub (which runs the full height from the pipeline row down to the
    #     resource boxes), so the connector line still cut through the "у" in
    #     "использует" at that y — same bug, different line.
    # (3) final fix: keep the y from (2) (clear of stop-arrow + text), but
    #     move the label OFF the connector's centerline — left-aligned,
    #     starting just right of where the connector passes — instead of
    #     centering across the full box width. No line crosses the label text
    #     at any point along its horizontal span now.
    conn_label_y = stop_y + 0.13
    conn_label_dx = 0.14  # clears the 0.06"-wide connector stub + margin
    text_box(s, x=tools_x + stage_w / 2 + conn_label_dx, y=conn_label_y, w=stage_w / 2, h=0.20,
             text="использует", size=8.5, italic=True, color=TEAL, align=PP_ALIGN.LEFT)
    text_box(s, x=mem_x + stage_w / 2 + conn_label_dx, y=conn_label_y, w=stage_w / 2, h=0.20,
             text="использует", size=8.5, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)

    text_box(s, x=0.55, y=6.62, w=12.25, h=0.35,
             text="ReAct (Reasoning + Acting) — Yao et al. 2022 (arXiv:2210.03629)",
             size=10.5, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    refs_of_slide(s, "s18")
    notes_with_sources(s, "s18")


def build_s19(p):
    """Agent at work — Fix-15 (Phase 12.6, 2026-05-13): split into 2 slides.

    s19 (this slide) = «Агент за работой: 200 PDF» — sequential steps with
    the tool used at each step. Visual: numbered steps top-to-bottom +
    tool-icon column on the right.
    Autonomy levels moved to NEW s19a.
    """
    s = blank(p)
    eyebrow_pill(s, "АГЕНТ")
    slide_title(s, "Агент за работой: 200 PDF — последовательность шагов.", size=26, y=0.85)
    # Left: case card (compact)
    cx_, cy_, cw_, ch_ = 0.55, 1.85, 4.0, 5.05
    ocean_box(s, cx_, cy_, cw_, ch_)
    text_box(s, x=cx_ + 0.25, y=cy_ + 0.20, w=cw_ - 0.5, h=0.4,
             text="Кейс — типовой для агента", size=13, bold=True, color=TEAL)
    text_box(s, x=cx_ + 0.25, y=cy_ + 0.65, w=cw_ - 0.5, h=1.5,
             text="200 PDF-отчётов.\nИзвлечь дату, контрагента, сумму.\nСобрать сводную таблицу.",
             size=14, bold=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=cx_ + 0.25, y=cy_ + 2.25, w=cw_ - 0.5, h=2.6,
             text="Не модель — нет специализированной модели «возьми 200 произвольных PDF».\n\n"
                  "Не чат — некомфортно копировать 200 файлов в окно.\n\n"
                  "Агент — естественный выбор.",
             size=11, color=DEEP, line_spacing=1.45, italic=True)
    # Right: 7 sequential steps with tool used per step
    sx, sy, sw, sh = cx_ + cw_ + 0.35, 1.85, 7.85, 5.05
    ocean_box(s, sx, sy, sw, sh)
    text_box(s, x=sx + 0.25, y=sy + 0.18, w=sw - 0.5, h=0.4,
             text="Что делает агент — пошагово, с указанием инструмента",
             size=13, bold=True, color=DEEP)
    # issue #153 QA fix #1 (Russification, presentation-critic P1): tool labels
    # were bare English phrases — translated to Russian, acronyms (OCR, PDF,
    # API, CSV, LLM) kept per §5.8 keep-list.
    steps = [
        # (num, action, tool)
        ("1", "Получить список файлов",          "файловая система"),
        ("2", "Открыть PDF #1",                   "чтение PDF"),
        ("3", "Извлечь текст",                    "извлечение текста (OCR / парсер)"),
        ("4", "Сделать сводку → vector DB",       "эмбеддинги + векторная БД"),
        ("5", "Найти ключевые поля",              "поиск + извлечение LLM"),
        ("6", "Записать строку в таблицу",        "запись в таблицу (Sheets API / CSV)"),
        ("7", "Цикл по всем 200 файлам",          "цикл оркестратора"),
    ]
    step_top = sy + 0.70
    step_h = 0.48
    step_gap = 0.10
    for i, (num, action, tool) in enumerate(steps):
        ry = step_top + i * (step_h + step_gap)
        is_loop = (i == len(steps) - 1)
        # Number badge
        bd = 0.40
        bx = sx + 0.30
        bcirc = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(bx), Inches(ry + (step_h - bd) / 2),
                                   Inches(bd), Inches(bd))
        bcirc.fill.solid()
        bcirc.fill.fore_color.rgb = GOLD if is_loop else MID
        bcirc.line.fill.background()
        disable_shadow(bcirc)
        text_box(s, x=bx, y=ry + (step_h - bd) / 2, w=bd, h=bd, text=num,
                 size=13, bold=True, color=DEEP if is_loop else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Action
        text_box(s, x=bx + bd + 0.20, y=ry + 0.08, w=4.20, h=step_h - 0.16,
                 text=action, size=12, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Tool — right-aligned in teal-tinted small block.
        # issue #153 QA fix #1: Russian tool phrases run longer than the EN
        # originals — shrink font for the longest labels to keep single-line.
        tx = sx + sw - 2.60
        filled_rect(s, tx, ry + 0.08, 2.30, step_h - 0.16, TEAL_TINT,
                    stroke=TEAL, stroke_pt=0.8, radius=True, radius_adj=0.30)
        tool_size = 8.5 if len(tool) > 28 else 10
        text_box(s, x=tx + 0.05, y=ry + 0.08, w=2.20, h=step_h - 0.16, text=tool,
                 size=tool_size, italic=True, color=TEAL,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.42, 12.25, 0.55,
                 "Агент = последовательность вызовов инструментов, оркестрируемая LLM. [1]",
                 size=13)
    refs_of_slide(s, "s19")
    notes_with_sources(s, "s19")


def build_s19a(p):
    """NEW slide — Fix-15 (Phase 12.6, 2026-05-13): autonomy levels.

    Pulled out of old s19. Shows 5 levels (Feng/McDonald/Zhang, 2025) +
    Human-in-the-loop / Human-on-the-loop / Human-out-of-the-loop / Override
    framing.
    """
    s = blank(p)
    eyebrow_pill(s, "АГЕНТ")
    slide_title(s, "Уровни автономии AI-агентов: дизайн-решение, не свойство модели.", size=24, y=0.85)
    # Left: 5 levels ladder
    lx, ly, lw, lh = 0.55, 1.85, 6.30, 5.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, x=lx + 0.25, y=ly + 0.20, w=lw - 0.5, h=0.4,
             text="5 уровней автономии (Feng / McDonald / Zhang, 2025) [1]",
             size=13, bold=True, color=DEEP)
    # issue #153 QA fix #1 (Russification, presentation-critic P1): level names
    # were bare English terms — Russian primary label + English gloss in
    # parens at first appearance, matching the s25 "Смещение (bias)" pattern.
    # issue #155 QA fix #192: role descriptions rewritten as explicit
    # "Пользователь: ..." statements so each card reads directly as "what the
    # user does at this level", not an abstract process description.
    levels = [
        ("1. Оператор (Operator)",     "Пользователь: одобряет каждое действие",  "Claude Code «approve each»", LIGHT),
        ("2. Соавтор (Collaborator)",  "Пользователь: работает наравне с агентом", "Cursor парное программирование", LIGHT),
        ("3. Консультант (Consultant)","Пользователь: ставит цель, правит план",  "Devin фиксит баг по тикету", MID),
        ("4. Утверждающий (Approver)", "Пользователь: утверждает на контрольных точках", "агент собирает PR, ждёт review", MID),
        ("5. Наблюдатель (Observer)",  "Пользователь: только получает результат", "AutoGPT на ночь", GOLD),
    ]
    rh = 0.78
    rt = ly + 0.80
    for i, (name, role, ex, color) in enumerate(levels):
        ry = rt + (4 - i) * rh  # bottom-up ladder визуально
        is_gold = (color == GOLD)
        filled_rect(s, lx + 0.25, ry, lw - 0.5, rh - 0.05, color,
                    stroke=DEEP if is_gold else None, stroke_pt=1.5 if is_gold else 0.0,
                    radius=True, radius_adj=0.15)
        # issue #153 QA fix #1: RU+EN-gloss names ("4. Утверждающий (Approver)")
        # run longer than the old bare EN names — name gets its own full-width
        # line at a smaller font, example line moved below it (was right-aligned
        # same line, no longer fits).
        text_box(s, x=lx + 0.40, y=ry + 0.05, w=lw - 0.8, h=0.28, text=name,
                 size=11.5, bold=True, color=DEEP if is_gold else WHITE)
        text_box(s, x=lx + 0.40, y=ry + 0.34, w=lw - 0.8, h=0.22, text=role,
                 size=9.5, italic=True, color=DEEP if is_gold else WHITE)
        text_box(s, x=lx + 0.40, y=ry + 0.54, w=lw - 0.8, h=0.20, text=ex,
                 size=9, color=DEEP if is_gold else WHITE, align=PP_ALIGN.RIGHT)
    # Right: 4 framings (in-the-loop / on-the-loop / out-of-the-loop / override)
    rx, ry_, rw, rh_ = lx + lw + 0.30, 1.85, SLIDE_W_IN - (lx + lw + 0.30) - 0.55, 5.05
    ocean_box(s, rx, ry_, rw, rh_)
    text_box(s, x=rx + 0.25, y=ry_ + 0.20, w=rw - 0.5, h=0.4,
             text="Где находится человек по отношению к циклу",
             size=13, bold=True, color=DEEP)
    # issue #153 QA fix #1 (Russification, presentation-critic P1): frame names
    # were bare English terms — Russian primary label + English gloss in
    # parens (matches s25 pattern). "Override" kept in parens as a concept
    # word, not left bare per brief.
    # issue #155 QA fix #193: order reversed (was in→on→out→Override, low
    # autonomy first) so rows line up with the left ladder, which has HIGH
    # autonomy (5. Наблюдатель) at the top: out-of-the-loop (≈ур.5) now top,
    # then on-the-loop (≈ур.3-4), then in-the-loop (≈ур.1-2); Override stays
    # last since it applies to any level, not a fixed ladder position.
    framings = [
        ("Человек вне цикла (Human-out-of-the-loop)", "Человек только смотрит итог (≈ ур. 5).",              GOLD),
        ("Человек над циклом (Human-on-the-loop)",    "Человек наблюдает, прерывает при отклонениях (≈ ур. 3-4).", MID),
        ("Человек в цикле (Human-in-the-loop)",       "Человек одобряет каждый шаг (≈ ур. 1-2).",            LIGHT),
        ("Режимы ручного перехвата (Override)",       "Любой уровень — с ручным перехватом по тревоге.",     TEAL),
    ]
    fy_top = ry_ + 0.75
    fh = 0.95
    fg = 0.10
    for i, (name, desc, color) in enumerate(framings):
        fy = fy_top + i * (fh + fg)
        ocean_box(s, rx + 0.20, fy, rw - 0.40, fh,
                  fill=GOLD_TINT if color == GOLD else WHITE,
                  stroke=color, stroke_pt=1.5)
        # Longer RU+EN-gloss names need a smaller font to stay single-line.
        name_size = 10.5 if len(name) > 30 else 12
        text_box(s, x=rx + 0.40, y=fy + 0.10, w=rw - 0.80, h=0.30,
                 text=name, size=name_size, bold=True, color=color)
        text_box(s, x=rx + 0.40, y=fy + 0.42, w=rw - 0.80, h=0.45,
                 text=desc, size=11, color=DEEP, line_spacing=1.30)
    # Bottom takeaway
    gold_callout(s, 0.55, 6.42, 12.25, 0.55,
                 "Уровень автономии — выбор продукта, не свойство модели. На семинарах будем выбирать осознанно.",
                 size=12)
    refs_of_slide(s, "s19a")
    notes_with_sources(s, "s19a")


def build_s20(p):
    """Applications product UX — Translate metrics + 6 logos grid."""
    s = blank(p)
    slide_title(s, "Приложение = AI, упакованный в продуктовый интерфейс.", size=28)
    # Top: Translate metrics
    # issue #153 QA fix #4 (P0, presentation-critic): the mixed 26pt/14pt
    # metrics line wrapped to 2 lines at the old height/box, overlapping the
    # fixed-position source caption below it. Fix: taller box (1.4→1.7),
    # dedicated single-line-height metrics row + caption pushed down to clear
    # the actual wrapped height, and font sizes trimmed slightly (26→24pt,
    # 14→13pt) to reduce wrap risk further. Also Russified "across" → "в".
    mt_x, mt_y, mt_w, mt_h = 0.55, 1.85, 12.25, 1.7
    ocean_box(s, mt_x, mt_y, mt_w, mt_h, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=mt_x + 0.3, y=mt_y + 0.18, w=mt_w - 0.6, h=0.35,
             text="Google Translate — масштаб 2026", size=14, bold=True, color=TEAL)
    text_runs(s, mt_x + 0.3, mt_y + 0.58, mt_w - 0.6, 0.75, [
        {"text": "1+ миллиард ", "size": 24, "bold": True, "color": DEEP},
        {"text": "пользователей в месяц  ·  ", "size": 13, "color": DEEP},
        {"text": "1 триллион ", "size": 24, "bold": True, "color": GOLD},
        {"text": "слов переведено / месяц", "size": 13, "color": DEEP},
    ], line_spacing=1.0)
    text_box(s, x=mt_x + 0.3, y=mt_y + mt_h - 0.38, w=mt_w - 0.6, h=0.3,
             text="в Google Translate, Search, Lens и Circle to Search (Google Blog, апрель 2026) [1]",
             size=10, italic=True, color=LIGHT)
    # 6 logo grid
    logos = [
        ("logo-googletranslate.png", "Google Translate", "нейронная трансляция"),
        ("logo-notion.png", "Notion AI", "GPT-4/Claude в кнопках"),
        ("logo-yandex.png", "ЯндексGPT в Поиске", "AI-краткий ответ"),
        ("logo-grammarly.png", "Grammarly", "NLP + LLM подсказки"),
        ("logo-yandex.png", "Яндекс.Карты", "ML маршрутизация"),
        ("logo-adobefirefly.png", "Adobe Firefly", "диффузия в Photoshop"),
    ]
    # grid_y nudged down 0.15" (3.55→3.70) to clear the taller metrics box
    # above (fix #4); cell_h trimmed 1.4→1.3 to keep the 2-row grid clear of
    # the gold_callout at y=6.55 below.
    grid_y = 3.70
    cell_w = 4.05
    cell_h = 1.3
    grid_x = 0.55
    cell_gap = 0.10
    for i, (icon, name, role) in enumerate(logos):
        col = i % 3
        row = i // 3
        x = grid_x + col * (cell_w + cell_gap)
        y = grid_y + row * (cell_h + cell_gap)
        ocean_box(s, x, y, cell_w, cell_h)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + 0.20, y=y + 0.30, w=0.80, h=0.80)
        text_box(s, x=x + 1.15, y=y + 0.25, w=cell_w - 1.3, h=0.4, text=name,
                 size=13, bold=True, color=DEEP)
        text_box(s, x=x + 1.15, y=y + 0.70, w=cell_w - 1.3, h=0.55, text=role,
                 size=10, italic=True, color=SLATE, line_spacing=1.30)
    gold_callout(s, 0.55, 6.47, 12.25, 0.50,
                 "AI как функция, а не продукт. Большинство студентов уже ежедневно пользуются всеми шестью. [2]",
                 size=12)
    refs_of_slide(s, "s20")
    notes_with_sources(s, "s20")


def build_s21(p):
    """Checklist 2 questions + quadrant — Fix-16 (Phase 12.6, 2026-05-13).

    New axes alignment:
    - Q1 (взаимодействие?) sits VERTICALLY on the left edge of the quadrant
      with «Да» beside the top row and «Нет» beside the bottom row.
    - Q2 (инструменты?) sits HORIZONTALLY below the quadrant with «Нет»
      under the left column and «Да» under the right column.
    Bottom takeaway («Подумайте 30 секунд…») removed.
    """
    s = blank(p)
    slide_title(s, "Чек-лист «Какой тип AI выбрать»: 2 вопроса + квадрант 2×2. [1] [2] [3]", size=24)
    # Quadrant area — large, centred. Shrunk vertically to leave room for Q2
    # markers + title + caption ABOVE slide bottom (7.5).
    quad_x, quad_y = 3.40, 1.65
    quad_w, quad_h = 7.50, 4.10
    filled_rect(s, quad_x, quad_y, quad_w, quad_h, WHITE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.04)
    # Cross
    filled_rect(s, quad_x + quad_w / 2 - 0.005, quad_y, 0.01, quad_h, SOFT_GREY)
    filled_rect(s, quad_x, quad_y + quad_h / 2 - 0.005, quad_w, 0.01, SOFT_GREY)

    # ─── Q1 axis (vertical, LEFT) — Fix-16 v4 ───
    # Layout in left column (no overlap with title):
    #   ВОПРОС 1 + question text  (CENTRE-vertically beside quadrant centre)
    #   ДА marker (beside top half centre)
    #   НЕТ marker (beside bottom half centre)
    # Markers placed FAR LEFT; question text wedged between markers.
    q1_x = 0.55
    q1_w = quad_x - q1_x - 0.30
    marker_w = 0.60
    marker_x = q1_x + 0.10
    text_x = marker_x + marker_w + 0.15
    text_w = q1_w - (marker_w + 0.25)
    # ДА marker — beside top half (centre y of top row)
    da_y_q1 = quad_y + quad_h / 4 - 0.18
    filled_rect(s, marker_x, da_y_q1, marker_w, 0.36,
                GOLD_TINT, stroke=GOLD, stroke_pt=1.0, radius=True, radius_adj=0.40)
    text_box(s, x=marker_x, y=da_y_q1, w=marker_w, h=0.36,
             text="ДА", size=12, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # НЕТ marker — beside bottom half
    net_y_q1 = quad_y + 3 * quad_h / 4 - 0.18
    filled_rect(s, marker_x, net_y_q1, marker_w, 0.36,
                WHITE, stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.40)
    text_box(s, x=marker_x, y=net_y_q1, w=marker_w, h=0.36,
             text="НЕТ", size=12, bold=True, color=LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Q1 title + question — vertically centred between markers
    centre_y = quad_y + quad_h / 2
    text_box(s, x=text_x, y=centre_y - 0.55, w=text_w, h=0.32,
             text="ВОПРОС 1",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=text_x, y=centre_y - 0.20, w=text_w, h=0.85,
             text="Нужно ли взаимодействие с пользователем?",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.LEFT, line_spacing=1.25)

    # ─── Q2 axis (horizontal, BOTTOM) — Fix-16 v2 ───
    # Markers right under quadrant (touching bottom edge), then Q2 title BELOW markers.
    q2_y = quad_y + quad_h + 0.10
    # Markers row — directly under quadrant columns
    marker_y = q2_y
    # Нет — under left column
    net_x = quad_x + quad_w / 4 - 0.30
    filled_rect(s, net_x, marker_y, 0.60, 0.36,
                WHITE, stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.40)
    text_box(s, x=net_x, y=marker_y, w=0.60, h=0.36,
             text="НЕТ", size=12, bold=True, color=LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Да — under right column
    da_x = quad_x + 3 * quad_w / 4 - 0.30
    filled_rect(s, da_x, marker_y, 0.60, 0.36,
                GOLD_TINT, stroke=GOLD, stroke_pt=1.0, radius=True, radius_adj=0.40)
    text_box(s, x=da_x, y=marker_y, w=0.60, h=0.36,
             text="ДА", size=12, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Q2 title — BELOW markers
    title_y = marker_y + 0.50
    text_box(s, x=quad_x, y=title_y, w=quad_w, h=0.32,
             text="ВОПРОС 2",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=quad_x, y=title_y + 0.32, w=quad_w, h=0.36,
             text="Нужна ли самостоятельная работа с инструментами?",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER)

    # ─── Cell labels with worked examples ───
    cells = [
        # (col, row, big, sub, color)
        (0, 0, "ЧАТ", "корп. чат для\nразбора ТЗ", LIGHT),
        (1, 0, "АГЕНТ", "200 PDF →\nсводная таблица", GOLD),
        (0, 1, "МОДЕЛЬ", "конвейерный\nдетектор", LIGHT),
        (1, 1, "ПРИЛОЖЕНИЕ", "ETL с AI-классификатором\n(автоматизация)", LIGHT),
    ]
    cw_ = quad_w / 2
    ch_ = quad_h / 2
    for col, row, big, sub, color in cells:
        cx = quad_x + col * cw_
        cy = quad_y + row * ch_
        is_gold = (color == GOLD)
        if is_gold:
            filled_rect(s, cx + 0.18, cy + 0.18, cw_ - 0.36, ch_ - 0.36, GOLD_TINT,
                        radius=True, radius_adj=0.08)
        text_box(s, x=cx + 0.20, y=cy + 0.30, w=cw_ - 0.4, h=0.65, text=big,
                 size=22, bold=True, color=DEEP if is_gold else color,
                 align=PP_ALIGN.CENTER, line_spacing=1.10)
        text_box(s, x=cx + 0.20, y=cy + 1.10, w=cw_ - 0.4, h=0.85, text=sub,
                 size=12, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.30)
    # Fix-16: bottom takeaway («Подумайте 30 секунд…») REMOVED.
    refs_of_slide(s, "s21")
    notes_with_sources(s, "s21")


def build_s22(p):
    """Section 4 divider — zoom-in state of unified nav template (Fix-17).

    Same 6-card grid as s02a/s10/s27; card 4 is gold-FILLED with white text.
    The 3 numbered "reasons" framing (data / quality / responsibility) was
    moved out of the slide visual to keep the navigation pattern uniform —
    the lecturer voices it via speaker notes (Лектору block in s22 md).
    """
    s = blank(p)
    # Fix-19: sub_marker removed — gold-filled active card is sole navigation indicator.
    nav_slide(s, here_idx=4,
              title="Раздел 4 · Границы AI — ваша зона ответственности",
              frame_phrase="Куда уходят данные · ошибки AI · «не умеет» — тоже ваше. [1] [2] [3]")
    refs_of_slide(s, "s22")
    notes_with_sources(s, "s22")


def build_s23(p):
    """Consumer vs enterprise — 2 columns + Samsung anchor + EU AI Act.

    issue #153 fix #15: bridge label added connecting from the section-4
    divider framing («Границы AI — ваша зона ответственности») to this
    slide's concrete topic (data destination).

    issue #155 fix #194: full visual overhaul — owner flagged the slide as
    "внезапно и неаккуратно". Fixes: (1) "ENTERPRISE / API" heading
    russified to match left column's russified heading, (2) single outer
    Ocean rounded box now wraps the whole composition (2 columns + bottom
    strip) instead of 4 disjoint floating islands, matching the s08/s12
    single-container pattern, (3) clearer top-to-bottom hierarchy:
    bridge-label -> title -> outer frame -> 2 columns -> bottom strip.
    """
    s = blank(p)
    slide_title(s, "Потребительские vs корпоративные тарифы — куда уходят ваши данные.", size=21, y=0.35, h=1.00)
    text_box(s, x=0.55, y=1.40, w=12.25, h=0.32,
             text="От общей зоны ответственности — к первому конкретному риску: данные.",
             size=13, italic=True, color=TEAL, align=PP_ALIGN.LEFT)
    # Single outer container wrapping the whole composition (fix #194)
    outer_x, outer_y = 0.55, 1.85
    outer_w, outer_h = SLIDE_W_IN - 2 * 0.55, 5.05   # -0.25 to open ref-list footer band
    ocean_box(s, outer_x, outer_y, outer_w, outer_h, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    pad = 0.25
    # Two columns (inside outer container)
    col_y, col_h = outer_y + pad, 3.30
    col_w = (outer_w - 2 * pad - 0.30) / 2
    # Left consumer
    cx_ = outer_x + pad
    ocean_box(s, cx_, col_y, col_w, col_h, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, x=cx_ + 0.25, y=col_y + 0.18, w=col_w - 0.5, h=0.40,
             text="ПОТРЕБИТЕЛЬСКИЕ ТАРИФЫ",
             size=14, bold=True, color=DEEP)
    text_box(s, x=cx_ + 0.25, y=col_y + 0.60, w=col_w - 0.5, h=0.45,
             text="данные → обучение по умолчанию",
             size=15, bold=True, color=DEEP, line_spacing=1.20)
    bullets_l = [
        "ChatGPT Free / Plus — обучение\nпо умолчанию",
        "Anthropic Claude (с сент. 2025) —\nпо согласию, 5 лет хранение",
        "Gemini Free — обучение +\nпроверка людьми, 3 года",
        "YandexGPT Free — стандартная\nполитика",
    ]
    for i, b in enumerate(bullets_l):
        text_box(s, x=cx_ + 0.30, y=col_y + 1.20 + i * 0.50, w=col_w - 0.6, h=0.50,
                 text=f"•  {b}", size=11, color=DEEP, line_spacing=1.12)
    # Right enterprise — heading russified (fix #194)
    ex_ = cx_ + col_w + 0.30
    ocean_box(s, ex_, col_y, col_w, col_h, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=ex_ + 0.25, y=col_y + 0.18, w=col_w - 0.5, h=0.40,
             text="КОРПОРАТИВНЫЕ ТАРИФЫ / API [2]",
             size=14, bold=True, color=TEAL)
    text_box(s, x=ex_ + 0.25, y=col_y + 0.60, w=col_w - 0.5, h=0.45,
             text="данные ≠ обучение",
             size=15, bold=True, color=DEEP, line_spacing=1.20)
    bullets_r = [
        "ChatGPT Enterprise / Business —\nбез обучения на данных",
        "OpenAI API (с марта 2023) —\nбез обучения на данных",
        "Anthropic for Business —\nнулевое хранение данных доступно",
        "Google Workspace / Vertex AI —\nбез обучения на данных",
    ]
    for i, b in enumerate(bullets_r):
        text_box(s, x=ex_ + 0.30, y=col_y + 1.20 + i * 0.50, w=col_w - 0.6, h=0.50,
                 text=f"•  {b}", size=11, color=DEEP, line_spacing=1.12)
    # Bottom strip (inside outer container): Samsung anchor + EU fines
    bot_y = col_y + col_h + 0.20
    bot_h = outer_y + outer_h - pad - bot_y
    s_w = (outer_w - 2 * pad) * 0.62
    s_x = outer_x + pad
    filled_rect(s, s_x, bot_y, s_w, bot_h, GOLD_TINT, stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.12)
    text_box(s, x=s_x + 0.20, y=bot_y + 0.10, w=s_w - 0.4, h=0.35,
             text="Samsung 2023 — канонический инцидент [1]", size=13, bold=True, color=DEEP)
    text_box(s, x=s_x + 0.20, y=bot_y + 0.48, w=s_w - 0.4, h=bot_h - 0.55,
             text="3 эпизода за месяц (март–апрель): код, транскрипт совещания, тестовые последовательности → попали в датасет OpenAI. Самсунг ввёл запрет внешнего GenAI.",
             size=11, color=DEEP, line_spacing=1.28)
    # EU AI Act
    eu_x = s_x + s_w + 0.25
    eu_w = outer_x + outer_w - pad - eu_x
    filled_rect(s, eu_x, bot_y, eu_w, bot_h, MID, radius=True, radius_adj=0.12)
    text_box(s, x=eu_x + 0.20, y=bot_y + 0.10, w=eu_w - 0.4, h=0.35,
             text="EU AI Act — штрафы [3]", size=13, bold=True, color=WHITE)
    text_box(s, x=eu_x + 0.20, y=bot_y + 0.48, w=eu_w - 0.4, h=0.33,
             text="до 15M € / 3% оборота", size=12, color=WHITE, bold=True)
    text_box(s, x=eu_x + 0.20, y=bot_y + 0.82, w=eu_w - 0.4, h=bot_h - 0.90,
             text="до 35M € / 7% — за запрещённые практики", size=11, color=GOLD, bold=True, line_spacing=1.15)
    refs_of_slide(s, "s23")
    notes_with_sources(s, "s23")


def build_s24(p):
    """Hallucinations — fake DOI prompt + Vectara HHEM range + AI knows all."""
    s = blank(p)
    slide_title(s, "Галлюцинации — неотъемлемое свойство AI. [1] [2]", size=28)
    # Left: prompt + 3 fake DOIs
    px, py, pw, ph = 0.55, 1.95, 7.5, 4.5
    ocean_box(s, px, py, pw, ph)
    text_box(s, x=px + 0.25, y=py + 0.20, w=pw - 0.5, h=0.45,
             text="Промпт", size=13, bold=True, color=TEAL)
    filled_rect(s, px + 0.25, py + 0.65, pw - 0.5, 0.7, SURFACE, stroke=SOFT_GREY,
                stroke_pt=1.0, radius=True, radius_adj=0.08)
    text_box(s, x=px + 0.40, y=py + 0.78, w=pw - 0.8, h=0.5,
             text='«Назови три статьи 2023-2024 по теме "сейсмостойкость подземных трубопроводов" с авторами, журналом и DOI».',
             size=11, italic=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=px + 0.25, y=py + 1.55, w=pw - 0.5, h=0.4,
             text="Ответ AI (3 фейк-ссылки):",
             size=13, bold=True, color=DEEP)
    fakes = [
        ("Petrov A., Smith J. (2023).", "Seismic Resilience of Small-Diameter Pipelines.", "DOI: 10.1016/j.engfailanal.2023.107214 ✗"),
        ("Ivanov K. et al. (2024).", "Underground Infrastructure Earthquake Response.", "DOI: 10.1080/15732479.2024.2218450 ✗"),
        ("Chen L., Brown R. (2023).", "Microscale Pipe Vibration Analysis.", "DOI: 10.1007/s11069-023-06122-1 ✗"),
    ]
    for i, (auth, title, doi) in enumerate(fakes):
        ry = py + 2.0 + i * 0.78
        text_box(s, x=px + 0.40, y=ry, w=pw - 0.8, h=0.30, text=auth,
                 size=10, bold=True, color=DEEP)
        text_box(s, x=px + 0.40, y=ry + 0.25, w=pw - 0.8, h=0.30, text=title,
                 size=10, italic=True, color=DEEP)
        text_box(s, x=px + 0.40, y=ry + 0.50, w=pw - 0.8, h=0.30, text=doi,
                 size=10, color=DEEP, bold=True)
    # Right: Vectara HHEM band
    rx, ry_, rw, rh = px + pw + 0.35, 1.95, 4.4, 3.2
    ocean_box(s, rx, ry_, rw, rh, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, x=rx + 0.25, y=ry_ + 0.20, w=rw - 0.5, h=0.4,
             text="Vectara HHEM (2025-26) [3]", size=13, bold=True, color=TEAL)
    text_box(s, x=rx + 0.25, y=ry_ + 0.65, w=rw - 0.5, h=0.4,
             text="процент галлюцинаций", size=11, italic=True, color=LIGHT)
    # Range bar
    text_box(s, x=rx + 0.25, y=ry_ + 1.25, w=rw - 0.5, h=0.45,
             text="< 1%", size=24, bold=True, color=TEAL)
    text_box(s, x=rx + 0.25, y=ry_ + 1.65, w=rw - 0.5, h=0.35,
             text="суммаризация (Gemini 2.0 Flash)",
             size=10, italic=True, color=SLATE)
    text_box(s, x=rx + 0.25, y=ry_ + 2.20, w=rw - 0.5, h=0.45,
             text="10–15%", size=24, bold=True, color=DEEP)
    text_box(s, x=rx + 0.25, y=ry_ + 2.60, w=rw - 0.5, h=0.35,
             text="reasoning (многошаговое)",
             size=10, italic=True, color=SLATE)
    # Anti-pattern callout below
    gold_callout(s, rx, ry_ + rh + 0.20, rw, 1.10,
                 "Анти-паттерн: «AI знает всё». Любой ответ AI по фактическому вопросу — гипотеза для проверки. [4]",
                 size=12)
    refs_of_slide(s, "s24")
    notes_with_sources(s, "s24")


def build_s25(p):
    """Смещение / лесть / дрейф распределения — issue #153 fix #16 (Russification).

    Card titles + assertion translated per chapter §4.4 / README §5.8 table:
    Bias → Смещение (bias); Sycophancy → Лесть (sycophancy);
    Distribution shift → Дрейф распределения (distribution shift).
    """
    s = blank(p)
    slide_title(s, "Смещение, лесть, дрейф распределения — три проявления одной природы.", size=24)
    cards = [
        ("Смещение (bias)", "lucide-scale-blue.png",
         "Модель повторяет перекосы датасета.",
         "Скрининг резюме обучен на исторических данных — дискриминирует, не «решая», а статистически."),
        ("Лесть (sycophancy)", "lucide-smartphone-blue.png",
         "Модель учится у разметки обратной связи поддакивать.",
         "Соглашается с явно неверным, чрезмерно хвалит — пользователь не замечает потери критики."),
        ("Дрейф распределения\n(distribution shift)", "lucide-trending-up-blue.png",
         "Данные периода — устаревают.",
         "Модель на коде 2023 в 2026 предложит устаревшую библиотеку без явного сбоя."),
    ]
    card_y = 1.85
    card_w = 4.05
    card_h = 3.4
    gap = 0.10
    start_x = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2.0
    colors = [LIGHT, MID, DEEP]
    for i, (name, icon, def_, ex) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        color = colors[i]
        ocean_box(s, x, card_y, card_w, card_h, stroke=color)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + 0.30, y=card_y + 0.30, w=0.65, h=0.65)
        title_size = 15 if "\n" in name else 20
        text_box(s, x=x + 1.10, y=card_y + 0.30, w=card_w - 1.3, h=0.6, text=name,
                 size=title_size, bold=True, color=color, line_spacing=1.10)
        text_box(s, x=x + 0.30, y=card_y + 1.20, w=card_w - 0.6, h=0.7, text=def_,
                 size=13, bold=True, color=DEEP, line_spacing=1.30)
        text_box(s, x=x + 0.30, y=card_y + 2.00, w=card_w - 0.6, h=1.2, text=ex,
                 size=11, italic=True, color=SLATE, line_spacing=1.40)
    # GPT-4o timeline
    tl_y = 5.50
    tl_h = 1.0
    ocean_box(s, 0.55, tl_y, 12.25, tl_h, fill=WHITE, stroke=GOLD, stroke_pt=2.0)
    text_box(s, x=0.75, y=tl_y + 0.10, w=11.85, h=0.4,
             text="GPT-4o: лесть (sycophancy) — апрель 2025 [1]", size=13, bold=True, color=DEEP)
    text_runs(s, 0.75, tl_y + 0.50, 11.85, 0.4, [
        {"text": "25 апр", "size": 14, "bold": True, "color": MID},
        {"text": " — релиз обновления   →   ", "size": 12, "color": DEEP},
        {"text": "28 апр", "size": 14, "bold": True, "color": MID},
        {"text": " — начало отката (Альтман в соцсети тем же вечером)   →   ", "size": 12, "color": DEEP},
        {"text": "29 апр", "size": 14, "bold": True, "color": MID},
        {"text": " — разбор причин", "size": 12, "color": DEEP},
    ])
    # Bottom takeaway
    text_box(s, x=0.55, y=6.57, w=12.25, h=0.4,
             text="Общая причина: модель отражает данные, на которых обучена. [2]",
             size=13, italic=True, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    refs_of_slide(s, "s25")
    notes_with_sources(s, "s25")


def build_s26(p):
    """4 speakers AGI table (renamed from old build_s27 in v3.1)."""
    s = blank(p)
    slide_title(s, "Прогнозы AGI — 4 спикера, 4 разных стимула. [1] [2]", size=26)
    # Table
    tx, ty, tw = 0.55, 1.95, 12.25
    rh_head = 0.5
    rh_row = 1.0
    # Header
    cols = [
        ("Спикер", 2.4),
        ("Аффилиация", 2.0),
        ("Прогноз AGI", 4.0),
        ("Материальный интерес", 3.85),
    ]
    # Header bg
    filled_rect(s, tx, ty, tw, rh_head, MID)
    cur_x = tx
    for label, w in cols:
        text_box(s, x=cur_x + 0.15, y=ty + 0.10, w=w - 0.3, h=rh_head - 0.2,
                 text=label, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w
    # Rows
    rows = [
        ("Sam Altman", "OpenAI",
         "«Знаем как построить AGI; начало superintelligence» (янв. 2026)",
         "$100B+ раунды; контекст IPO"),
        ("Dario Amodei", "Anthropic",
         "«AGI через 2–3 года; нобелевский уровень за 2 года» (Давос 2026)",
         "Конкуренция с OpenAI; раунд 2026"),
        ("Demis Hassabis", "Google\nDeepMind",
         "«AGI к 2029–2030 (3–4 года); окно сузилось за 2026 год» (Axios/Google I/O, май 2026)",
         "Лидер community; больше доверия при осторожной позиции"),
        ("Yann LeCun", "AMI Labs\n(экс-Meta)",
         "«LLM не приведут к AGI; нужны world models, JEPA»",
         "Раунд $1B (март 2026) на альтернативный путь"),
    ]
    for i, (sp, af, pr, st) in enumerate(rows):
        rt = ty + rh_head + i * rh_row
        bg = SURFACE if i % 2 == 0 else WHITE
        filled_rect(s, tx, rt, tw, rh_row, bg, stroke=SOFT_GREY, stroke_pt=0.5)
        cur_x = tx
        for j, (text, w) in enumerate(zip([sp, af, pr, st], [c[1] for c in cols])):
            is_speaker = (j == 0)
            text_box(s, x=cur_x + 0.15, y=rt + 0.15, w=w - 0.3, h=rh_row - 0.30,
                     text=text, size=11.5 if is_speaker else 10.5,
                     bold=is_speaker, color=DEEP if is_speaker else SLATE,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.30)
            cur_x += w
    # issue #155 fix #196: closing callout removed — slide ends with the
    # table, conclusion stays only in speaker notes.
    refs_of_slide(s, "s26")
    notes_with_sources(s, "s26")


def build_s27(p):
    """Section 5 divider — zoom-in state of unified nav template (Fix-17).

    Same 6-card grid as s02a/s10/s22; card 5 is gold-FILLED with white text.
    """
    s = blank(p)
    nav_slide(s, here_idx=5,
              title="Раздел 5 · Что забрать домой",
              frame_phrase="Резюме · карта семестра · тизер лекции 2.",
              sub_marker="↓ Финал — раздел 5 из 5")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """Summary — 3 takeaway cards. issue #153 fix #18: homework callout removed
    entirely (seminar assignment lives in the seminar, not the lecture)."""
    s = blank(p)
    slide_title(s, "Что мы прошли: три главных вывода.", size=28)
    takeaways = [
        ("AI — спектр, не монолит", "Тип задачи × модальность × тип реализации.\nГрамотное обсуждение начинается с явной классификации."),
        ("Выбор типа AI — навык", "2 диагностических вопроса + квадрант 2×2.\nИнструмент, который вы применяете на семинарах."),
        ("Целеполагание у человека", "Все классы ошибок требуют человеческого контура.\nГраница «AI / не-AI» — ваша инженерная зона."),
    ]
    card_y = 2.6
    card_w = 4.05
    card_h = 3.6
    gap = 0.10
    start_x = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2.0
    colors = [LIGHT, MID, DEEP]
    for i, (head, body) in enumerate(takeaways):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h, stroke=colors[i])
        text_box(s, x=x + 0.20, y=card_y + 0.20, w=0.7, h=0.85, text=str(i + 1),
                 size=44, bold=True, color=colors[i])
        text_box(s, x=x + 1.0, y=card_y + 0.30, w=card_w - 1.2, h=0.85, text=head,
                 size=15, bold=True, color=DEEP, line_spacing=1.20)
        text_box(s, x=x + 0.30, y=card_y + 1.40, w=card_w - 0.6, h=1.5, text=body,
                 size=11.5, color=DEEP, line_spacing=1.40)
    gold_callout(s, 0.55, 6.55, 12.25, 0.60,
                 "Главный вопрос лекции: где AI работает, где — нет, и как это понять?",
                 size=15)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """Course roadmap — issue #153 fix #19: FULL REDESIGN, 4-module structure.

    Canon (chapter.md v3.3 §5.1, issue #153, verified via GitHub REST API):
    - M1 = lectures 1.1-1.6 (6 lectures): theoretical-methodological foundation
      + industries closest to everyday engineering/product practice.
    - M2 = lectures 2.1-2.5 (5 lectures): engineering design, heavy industry,
      manufacturing, agriculture.
    - M3 = lectures 3.1-3.6 (6 lectures): infrastructure, science, resource
      extraction, final synthesis.
    - M4 = exam (30h incl. prep).
    РК1/РК2/РК3 — on completion of each of the first three modules (not tied
    to specific seminar numbers, which are out of scope per issue #153).
    Isolation: catalog/manifests/lectures.yaml and RPD NOT touched (issue #154).
    """
    # issue #153 QA fix #5 (P2, presentation-critic + student-simulator both
    # flagged): title "17 лекций × 4 модуля" read as if all 4 were the same
    # kind of block, when Module 4 is the final exam, not a content module.
    # Fix: title reworded to name 3 content modules + exam separately; the
    # 4th column's header label changed from "Модуль 4" to "Экзамен" (primary
    # label distinguishes it visually/textually) — column count/data/width
    # unchanged (structure stays as approved in issue #153).
    s = blank(p)
    slide_title(s, "Карта семестра: 17 лекций, 3 модуля + экзамен.", size=28)
    modules = [
        ("Модуль 1", "Теор.-метод. основы\n+ знакомые индустрии", LIGHT, [
            ("1.1", "1.1 Введение", True),
            ("1.2", "1.2 Архитектура AI", False),
            ("1.3", "1.3 Агенты, RAG, API", False),
            ("1.4", "1.4 Разработка ПО", False),
            ("1.5", "1.5 Финансы / ритейл", False),
            ("1.6", "1.6 Креативные ◆РК1", False),
        ]),
        ("Модуль 2", "Инж. проектирование,\nтяж. промышленность, АПК", MID, [
            ("2.1", "2.1 CAD/CAM", False),
            ("2.2", "2.2 Авиакосмос / ОПК", False),
            ("2.3", "2.3 Производство", False),
            ("2.4", "2.4 Цифровые двойники", False),
            ("2.5", "2.5 Сельское хоз-во ◆РК2", False),
        ]),
        ("Модуль 3", "Инфраструктура, наука,\nдобыча, синтез", DEEP, [
            ("3.1", "3.1 Логистика", False),
            ("3.2", "3.2 Телеком / cybersec", False),
            ("3.3", "3.3 Наука", False),
            ("3.4", "3.4 Нефтегаз", False),
            ("3.5", "3.5 Медицина", False),
            ("3.6", "3.6 Синтез ◆РК3", False),
        ]),
        # issue #153 QA fix #5 follow-up: original subheader "Итоговая
        # аттестация\n(не модуль с лекциями)" wrapped to 3 lines in the
        # narrow column and got clipped/overlapped the row below — shortened
        # to fit the same 2-line budget as the other 3 modules' subheaders.
        ("Экзамен", "Итоговая\nаттестация", TEAL, [
            ("Экз.", "Экзамен", False),
        ]),
    ]
    mod_y = 1.75
    mod_h = 4.55
    # Module width by weighted units (not raw lecture count) — M4 has only
    # 1 "row" but needs enough width for «Модуль 4» + «Экзамен» header text
    # to render without overlap (fix #19 bugfix: was raw n=1 unit, too narrow).
    weights = [len(lectures) for _, _, _, lectures in modules]
    weights[-1] = 2.2  # М4 minimum width guarantee
    total_units = sum(weights)
    bar_x = 0.55
    bar_w = SLIDE_W_IN - 2 * 0.55
    unit_w = bar_w / total_units
    cur_x = bar_x
    for (label, sub, color, lectures), w_units in zip(modules, weights):
        n = len(lectures)
        m_w = w_units * unit_w
        ocean_box(s, cur_x, mod_y, m_w - 0.05, mod_h, fill=WHITE, stroke=color, stroke_pt=2.0)
        # issue #155 QA fix P2-9 (same sliver root cause as s02a — see the
        # matching comment there): header strip radius must match ocean_box's
        # absolute ~12pt corner radius, computed against the strip's own
        # (shorter) height, not a fixed fractional radius_adj.
        head_h = 0.85
        head_radius_adj = max(0.04, min(0.35, (12.0 / 72.0) / (head_h / 2.0)))
        filled_rect(s, cur_x, mod_y, m_w - 0.05, head_h, color, radius=True, radius_adj=head_radius_adj)
        text_box(s, x=cur_x, y=mod_y + 0.08, w=m_w - 0.05, h=0.35, text=label,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text_box(s, x=cur_x + 0.10, y=mod_y + 0.40, w=m_w - 0.25, h=0.5, text=sub,
                 size=9.5, italic=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.15)
        for j, (lec_num, lec_label, is_now) in enumerate(lectures):
            ly = mod_y + 1.05 + j * 0.55
            text_box(s, x=cur_x + 0.15, y=ly, w=m_w - 0.3, h=0.50, text=lec_label,
                     size=11 if is_now else 10,
                     bold=True, color=DEEP,
                     line_spacing=1.20)
        cur_x += m_w
    text_box(s, x=0.55, y=6.55, w=12.25, h=0.35,
             text="◆ — рубежные контроли РК1/РК2/РК3, по завершении каждого из первых трёх модулей.",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s29"))


def build_s29a(p):
    """Grading formula strip — issue #153 fix #19 (new slide next to s29).

    Short one-formula slide: 100 = 10 (attendance) + 30 (exam) + 3×20 (РК).
    Kept minimal per brief — not a full content slide.
    """
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # issue #155 fix #198: explicit slide title added — slide previously
    # opened straight on the formula with no heading.
    slide_title(s, "Оценка за семестр", size=28, align=PP_ALIGN.CENTER)
    text_runs(s, 0.4, 3.15, 12.53, 1.0, [
        {"text": "100", "size": 52, "bold": True, "color": GOLD},
        {"text": "  =  10 ", "size": 30, "bold": True, "color": DEEP},
        {"text": "(посещаемость)", "size": 14, "italic": True, "color": SLATE},
        {"text": "  +  30 ", "size": 30, "bold": True, "color": DEEP},
        {"text": "(экзамен)", "size": 14, "italic": True, "color": SLATE},
        {"text": "  +  3×20 ", "size": 30, "bold": True, "color": DEEP},
        {"text": "(РК1/РК2/РК3)", "size": 14, "italic": True, "color": SLATE},
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.8, y=4.55, w=11.7, h=0.6,
             text="Рубежные контроли — по завершении каждого из первых трёх модулей.",
             size=16, italic=True, color=MID, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s29a"))


def build_s30(p):
    """Lecture 2 teaser — Fix-19 (Phase 12.6, 2026-05-13).

    Removed YOLO callback frame (visual repetition with s01). Now: title +
    4 concept cards (full slide width) + 1-phrase frame at the bottom.
    """
    s = blank(p)
    slide_title(s, "Лекция 2: «Как работают современные большие модели».", size=28)
    # 4 concept cards 2×2, centred on full slide
    grid_w = 11.50
    grid_x = (SLIDE_W_IN - grid_w) / 2.0
    grid_y = 2.05
    grid_h = 4.30
    sub_w = (grid_w - 0.30) / 2
    sub_h = (grid_h - 0.30) / 2
    concepts = [
        ("lucide-file-text-blue.png", "Токены",
            "единицы, на которые модель режет текст"),
        ("lucide-network-blue.png",   "Эмбеддинги",
            "(векторные представления) — адреса в смысловом пространстве"),
        ("lucide-eye-blue.png",       "Внимание",
            "(attention) — на какие части входа смотреть"),
        ("lucide-zap-blue.png",       "Температура",
            "случайность выбора следующего токена"),
    ]
    for i, (icon, name, sub) in enumerate(concepts):
        col = i % 2
        row = i // 2
        x = grid_x + col * (sub_w + 0.30)
        y = grid_y + row * (sub_h + 0.30)
        ocean_box(s, x, y, sub_w, sub_h)
        if (ASSETS / "icons" / icon).exists():
            add_image(s, ASSETS / "icons" / icon, x=x + 0.35, y=y + 0.35, w=0.85, h=0.85)
        text_box(s, x=x + 1.40, y=y + 0.40, w=sub_w - 1.55, h=0.55, text=name,
                 size=22, bold=True, color=MID)
        text_box(s, x=x + 0.35, y=y + 1.40, w=sub_w - 0.7, h=sub_h - 1.6, text=sub,
                 size=13, italic=True, color=DEEP, line_spacing=1.40)
    # Bottom 1-phrase frame
    text_box(s, x=0.55, y=6.65, w=12.25, h=0.45,
             text="Эти 4 концепта объясняют поведение всех современных LLM — от ChatGPT до DeepSeek.",
             size=14, italic=True, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """Вопросы? — minimal. issue #153 fix #21: renamed from «Q&A» (title text
    only, no redesign). Font size reduced 140pt → 96pt to fit the longer word."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=2.35, w=12.25, h=1.9, text="Вопросы?",
             size=96, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.0)
    text_box(s, x=0.55, y=5.4, w=12.25, h=0.7, text="Спасибо",
             size=36, color=MID, align=PP_ALIGN.CENTER, italic=True)
    # issue #155 fix #199: contact placeholder removed entirely (was
    # "контакты лектора — заполняется перед лекцией") — real contacts to
    # be added in a separate point pass later, not left as a stub.
    speaker_notes(s, load_notes("s31"))


# ============================================================
# Main
# ============================================================
BUILDERS = [
    build_s01,
    # issue #153 fix #2: s00a (welcome) + s00b (course hook, ex-s05b) inserted
    # between s01 (ice-breaker demo) and s02 (cover).
    build_s00a, build_s00b,
    build_s02, build_s02a,
    # issue #153 fix #1: build_s03/build_s04 (icebreaker poll) DELETED —
    # poll moves to seminar 1, out of scope for lec-01 slides.
    # issue #153 fix #2: build_s05b DELETED — content moved to build_s00b
    # (reworded role: hook before cover, not "course frame after instructor").
    build_s05a,
    # issue #155 fix #177: build_s05c (section 1 divider) NEW, inserted
    # between build_s05a and build_s06 — section-divider audit found no
    # divider before раздел 1.
    build_s05c,
    build_s06,
    # issue #153 fix #4: build_s06a (McCulloch-Pitts 1943 fact-bridge) NEW,
    # inserted between build_s06 and build_s07.
    build_s06a,
    build_s07,
    # issue #155 fix #177: build_s07a (section 2 divider) NEW, inserted
    # between build_s07 and build_s08 — section-divider audit found no
    # divider before раздел 2.
    build_s07a,
    build_s08, build_s09,
    build_s10, build_s11, build_s12, build_s13,
    # Fix-17 (2026-05-13): build_s14 (mini-divider «Разберём подробнее») deleted.
    # Reason: paraphrased s10 framing, used 4-type icons inconsistent with the
    # 5-section navigation grammar of s02a/s10/s22/s27. Verbal transition moved
    # to s13 «Лектору» speaker notes. Pacing recovered: -0.5 min from active.
    build_s15, build_s16, build_s17, build_s18, build_s19, build_s19a, build_s20, build_s21,
    build_s22, build_s23, build_s24, build_s25, build_s26,
    # v3.1: removed build_s26-old (ARC-AGI) and build_s28-old (Pearl);
    # added NEW build_s27 (section 5 divider); renumbered s27→s26, s29→s28, s30→s29, s31→s30, s32→s31.
    build_s27, build_s28,
    build_s29,
    # issue #153 fix #19: build_s29a (grading formula strip) NEW, after build_s29.
    build_s29a,
    build_s30, build_s31,
]


def main():
    p = setup_pres()
    for build in BUILDERS:
        build(p)
    # Stamp a muted «N / TOTAL» page number on every slide (bottom-right).
    total = len(p.slides)
    for i, slide in enumerate(p.slides, start=1):
        page_number(slide, i, total)
    p.save(str(OUT))
    print(f"Saved {OUT} with {len(BUILDERS)} slides (page numbers 1..{total}).")


if __name__ == "__main__":
    main()
