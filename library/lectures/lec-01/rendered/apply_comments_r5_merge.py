"""
apply_comments_r5_merge.py — re-applies editor comments #249-#256 ON TOP of the
reference-annotated lec-01.pptx that lives in origin/main (added by the publishing
session: inline [N] citation markers, 'N / 36' slide-number footers, expanded notes).

This merges my 8 edits WITHOUT losing the numbers/references:
  - text/shape finders tolerate the appended [N] markers
  - after reordering, ALL 'N / 36' footers are renumbered and the new chat slide
    gets its own footer
  - speaker notes on existing slides are never touched (kept as-is)

Base = current working-tree lec-01.pptx (must be the origin/main version).
"""
import copy
import re
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import build_lec01 as B
from build_lec01 import (
    text_box, ocean_box, filled_rect, chip, add_image, set_slide_bg,
    speaker_notes, disable_shadow,
    DEEP, MID, LIGHT, TEAL, SURFACE, WHITE, GOLD, SLATE, GOLD_TINT, SOFT_GREY,
)

HERE = Path(__file__).resolve().parent
PPTX = HERE / "lec-01.pptx"
SEM01 = HERE / "../../../seminars/sem-01/rendered/sem-01.pptx"
QR = HERE / "assets/max-qr.png"

prs = Presentation(str(PPTX))
S = prs.slides
FOOTER_RE = re.compile(r'\d+\s*/\s*\d+')


def emu_in(v):
    return Emu(v).inches if v is not None else 0.0


def set_lines(shp, lines):
    paras = shp.text_frame.paragraphs
    for i, line in enumerate(lines):
        if i < len(paras) and paras[i].runs:
            paras[i].runs[0].text = line
            for r in paras[i].runs[1:]:
                r.text = ''


def set_full_text(shp, text):
    para = shp.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r.text = ''
    else:
        para.text = text


def shape_by_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


def is_footer(sh):
    return sh.has_text_frame and FOOTER_RE.fullmatch(sh.text_frame.text.strip()) is not None


# ============================================================
# #249 — slide 10: reflow "Открытия" row 3 -> 4 items (add 1956 Dartmouth)
# ============================================================
def fix_s10():
    s = S[9]
    def find(txt):
        return shape_by_text(s, txt)
    t_turing = find('Тьюринг'); y_1950 = find('1950')
    t_eliza = find('ELIZA'); y_1966 = find('1966')
    t_expert = find('Экспертные системы'); y_1980 = find('1980-е')
    ticks = [sh for sh in s.shapes if sh.shape_type == 1 and abs(emu_in(sh.top) - 2.69) < 0.05
             and emu_in(sh.width) < 0.2]
    ticks.sort(key=lambda sh: emu_in(sh.left))

    W = 2.72
    xs = [1.00, 3.85, 6.70, 9.55]
    centers = [x + W / 2 for x in xs]
    for sh, col in [(t_turing, 0), (t_eliza, 2), (t_expert, 3)]:
        sh.left = Inches(xs[col]); sh.width = Inches(W); sh.top = Inches(2.05); sh.height = Inches(0.72)
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)
        sh.text_frame.word_wrap = True
    for sh, col in [(y_1950, 0), (y_1966, 2), (y_1980, 3)]:
        sh.left = Inches(xs[col]); sh.width = Inches(W)
    for sh, col in zip(ticks, [0, 2, 3]):
        sh.left = Inches(centers[col] - emu_in(sh.width) / 2)

    new_title_el = copy.deepcopy(t_turing._element)
    s.shapes._spTree.append(new_title_el)
    nt = s.shapes[-1]
    set_lines(nt, ['Дартмутский семинар — рождение термина «AI»'])
    nt.left = Inches(xs[1]); nt.width = Inches(W); nt.top = Inches(2.05); nt.height = Inches(0.72)
    for p in nt.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11)
    nt.text_frame.word_wrap = True

    new_year_el = copy.deepcopy(y_1950._element)
    s.shapes._spTree.append(new_year_el)
    ny = s.shapes[-1]; set_lines(ny, ['1956']); ny.left = Inches(xs[1]); ny.width = Inches(W)

    if ticks:
        new_tick_el = copy.deepcopy(ticks[0]._element)
        s.shapes._spTree.append(new_tick_el)
        ntk = s.shapes[-1]; ntk.left = Inches(centers[1] - emu_in(ntk.width) / 2)
    print('  #249 s10: 4-column Открытия row (added 1956 Дартмут)')


# ============================================================
# #250 — slide 12: refresh ChatGPT WAU stat (preserve [N] ref)
# ============================================================
def fix_s12():
    s = S[11]
    changed = []
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == '900M':
            set_full_text(sh, '1B'); changed.append('900M->1B')
        elif t.startswith('ChatGPT, февраль 2026'):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if 'февраль' in r.text:
                        r.text = r.text.replace('февраль', 'август'); changed.append('date(kept ref)')
    head = shape_by_text(s, '900M пользователей')
    if head:
        for p in head.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('900M пользователей', '~1 млрд пользователей')
        changed.append('headline')
    print('  #250 s12:', changed)


# ============================================================
# #253 — slide 22: visualize inner (1 file) + outer (x200) loops
# ============================================================
def _dash(line):
    ln = line._get_or_add_ln()
    d = ln.find(qn('a:prstDash'))
    if d is None:
        d = ln.makeelement(qn('a:prstDash'), {}); ln.append(d)
    d.set('val', 'dash')


def fix_s22():
    s = S[21]
    x, y, w, h = 5.03, 3.05, 7.66, 2.85
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background(); box.line.color.rgb = TEAL; box.line.width = Pt(1.6)
    _dash(box.line); disable_shadow(box)
    try:
        box.adjustments[0] = 0.05
    except Exception:
        pass
    lab = s.shapes.add_textbox(Inches(11.75), Inches(3.95), Inches(2.6), Inches(0.4))
    lab.rotation = 270
    ltf = lab.text_frame; ltf.word_wrap = False
    lp = ltf.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run(); lr.text = 'цикл: обработка одного файла'
    lr.font.size = Pt(11); lr.font.bold = True; lr.font.color.rgb = TEAL

    ax = 4.965
    bar = filled_rect(s, ax, 3.30, 0.055, 3.00, TEAL); disable_shadow(bar)
    stub = filled_rect(s, ax, 6.25, 0.30, 0.055, TEAL); disable_shadow(stub)
    head = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ax - 0.06), Inches(3.10), Inches(0.18), Inches(0.22))
    head.fill.solid(); head.fill.fore_color.rgb = TEAL; head.line.fill.background()
    disable_shadow(head)
    chip(s, 5.20, 6.46, 2.55, 0.32, '↻  повторить × 200 файлов', fill=GOLD, color=DEEP, size=11, bold=True)
    print('  #253 s22: inner dashed box + outer return arrow + labels')


# ============================================================
# #255 — slide 34: grading formula (runs unchanged in main version)
# ============================================================
def fix_s34():
    s = S[33]
    sh = shape_by_text(s, 'посещаемость')
    runs = sh.text_frame.paragraphs[0].runs
    runs[1].text = '  =  40 '
    runs[2].text = '(экзамен)'
    runs[3].text = ''
    runs[4].text = ''
    print('  #255 s34: grading formula -> 40 exam, no attendance')


# ============================================================
# #256 support — cross-file copy of the sem-01 MAX chat slide
# ============================================================
def copy_max_slide():
    src_prs = Presentation(str(SEM01))
    src = src_prs.slides[4]
    dest = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    src_csld = src._element.find(qn('p:cSld'))
    dst_csld = dest._element.find(qn('p:cSld'))
    bg = src_csld.find(qn('p:bg'))
    if bg is not None:
        dst_csld.insert(0, copy.deepcopy(bg))
    else:
        set_slide_bg(dest, WHITE)
    for shp in src.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    for pic in dest.shapes._spTree.iter(qn('p:pic')):
        blip = pic.find('.//' + qn('a:blip'))
        if blip is None:
            continue
        rId = blip.get(qn('r:embed'))
        if not rId:
            continue
        src_img = src.part.related_part(rId)
        img_part, new_rId = dest.part.get_or_add_image_part(BytesIO(src_img.blob))
        blip.set(qn('r:embed'), new_rId)
    speaker_notes(dest, MAX_NOTE)
    return list(prs.slides._sldIdLst)[-1]


MAX_NOTE = (
    "Всё общение по курсу — в чате в мессенджере MAX. Отсканируйте QR-код камерой телефона "
    "или перейдите по ссылке max.ru/join/sHoHlhI4jvW_swcq5ZLowvz6G94m0IMhAwngxuxsEpI "
    "и вступите в чат до первой лекции.\n\n"
    "Что там будет. Во-первых, объявления: расписание, любые изменения, важные новости — "
    "всё публикуется в чате, отдельно письма ждать не нужно. Во-вторых, материалы: слайды "
    "лекций и семинаров я выкладываю туда по мере проведения занятий, так что всё собрано в "
    "одном месте. В-третьих, вопросы: если что-то непонятно между парами — спрашивайте в "
    "чате, не обязательно копить до семинара.\n\n"
    "Пожалуйста, вступите сегодня — так вы точно не пропустите первое объявление."
)


def add_qr_to_qa():
    s = S[35]
    cx, cy, cw, ch = 9.7, 4.25, 3.1, 2.55
    ocean_box(s, cx, cy, cw, ch, fill=WHITE, stroke=LIGHT, stroke_pt=1.6)
    qsz = 1.55
    add_image(s, str(QR), cx + (cw - qsz) / 2, cy + 0.22, w=qsz, h=qsz)
    text_box(s, cx + 0.15, cy + 1.86, cw - 0.3, 0.32, 'Чат курса в MAX',
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, cx + 0.12, cy + 2.16, cw - 0.24, 0.4,
             'max.ru/join/sHoHlhI4jvW…', size=9.5, color=MID, align=PP_ALIGN.CENTER)
    print('  #256 s36: MAX QR added to Q&A slide')


# ============================================================
# Renumber all 'N / M' footers to final positions + add one to chat slide
# ============================================================
def renumber_footers():
    slides = list(prs.slides)
    n = len(slides)
    template = None
    for s in slides:
        for sh in s.shapes:
            if is_footer(sh):
                template = sh
                break
        if template is not None:
            break
    added = 0
    for i, s in enumerate(slides, 1):
        found = next((sh for sh in s.shapes if is_footer(sh)), None)
        if found is not None:
            set_full_text(found, f'{i} / {n}')
        elif template is not None:
            el = copy.deepcopy(template._element)
            s.shapes._spTree.append(el)
            newsh = s.shapes[-1]
            set_full_text(newsh, f'{i} / {n}')
            added += 1
    print(f'  footers renumbered 1..{n}; added {added} new footer(s)')


# ============================================================
# Orchestrate
# ============================================================
print('Content edits (on reference-annotated base):')
fix_s10()
fix_s12()
fix_s22()
fix_s34()
add_qr_to_qa()

print('Structural edits:')
lst = prs.slides._sldIdLst
ids = list(lst)
assert len(ids) == 36, f'expected 36, got {len(ids)}'

chat_id = copy_max_slide()
print('  #256: MAX chat slide copied from sem-01')

head_intro = ids[0:6]
moved = [ids[32], ids[33]]
tail = [ids[34], ids[35]]
middle = ids[6:32]
middle[8], middle[9] = middle[9], middle[8]      # #251 swap s15<->s16
del_el = ids[19]                                  # #252 delete s20
middle = [e for e in middle if e is not del_el]

desired = head_intro + moved + [chat_id] + middle + tail
assert len(desired) == 36, f'desired {len(desired)}'

del_rId = del_el.get(qn('r:id'))
prs.part.drop_rel(del_rId)
print('  #252: slide 20 (ЧАТ кейс) deleted')
print('  #251: slides 15<->16 swapped')
print('  #256: map+grading moved to front, chat inserted after grading')

for el in list(lst):
    lst.remove(el)
for el in desired:
    lst.append(el)

renumber_footers()

prs.save(str(PPTX))
print(f'Saved {PPTX} — total slides now {len(list(prs.slides))}')
