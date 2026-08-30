"""
Full 33-slide build of Лекции 5 «AI в финансовом секторе и ритейле».

Source-of-truth: deck.yaml + deck-part2.yaml (split >600 строк, loader reads
both) + chapter v2 finalized (3 части, ~22650 слов) + slides/*.md (33 файла:
32 LOCKED s01–s32 + s04a divider Раздела 1, readable speaker notes 150–300).

Issue #100 · worktree /tmp/lec-05-wt (branch phase-1-plan)

Palette LOCKED: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide. Motif «Ocean rounded box»
(radius 12, surface #F4F7FA, stroke #1C7293 1.5pt) на каждом content-слайде.

Canvas: 13.333" × 7.5" (16:9, [#55-1] patch). roadmap-bar ТОЛЬКО на
divider'ах (s04a/s10/s15/s20/s25) + cover (s02) — Л2-урок #40.

Render-style эталон: build_lec04.py (та же палитра/motif/типографика/
divider-шаблон/плотность). Build: python3 build_lec05.py → lec-05.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED) ===
DEEP = RGBColor(0x21, 0x29, 0x5C)
MID = RGBColor(0x06, 0x5A, 0x82)
LIGHT = RGBColor(0x1C, 0x72, 0x93)
TEAL = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF0, 0xAB, 0x00)
SLATE = RGBColor(0x5B, 0x66, 0x78)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFD, 0xF3, 0xDC)
TEAL_TINT = RGBColor(0xE4, 0xF1, 0xF2)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path("/home/harness/harness-projects/256/.worktrees/folder-288/publish-8a63bf98/library/lectures/lec-05")
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
CHARTS = ASSETS / "charts"
DIAGRAMS = ROOT / "rendered/assets-en/diagrams"   # EN diagrams (no Cyrillic)
SLIDES_DIR = ROOT / "slides-en"
OUT = ROOT / "rendered/lec-05-en.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers (architecture mirrors lec-04 build — proven)
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    pgr = tf.paragraphs[0]
    pgr.alignment = align
    pgr.line_spacing = line_spacing
    r = pgr.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    pgr = tf.paragraphs[0]
    pgr.alignment = align
    pgr.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            pgr = tf.add_paragraph()
            pgr.alignment = cfg.get("align", align)
            pgr.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before") is not None:
                pgr.space_before = Pt(cfg["space_before"])
        r = pgr.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.035, min(0.22,
                             (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, fill, *, stroke=None, stroke_pt=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    pgr = tf.paragraphs[0]
    pgr.alignment = PP_ALIGN.CENTER
    r = pgr.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def connector(slide, x1, y1, x2, y2, color=LIGHT, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = etree.SubElement(
            ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                "prstDash")
        pd.set("val", dash)
    return cn


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    """[#73-render-1] aspect-safe: pass only the constraining dimension."""
    path = Path(path)
    if not path.exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path)
            iw, ih = img.size
            img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        ir = iw / ih
        br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x),
                                     Inches(y + (h - ah) / 2), width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2),
                                     Inches(y), height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.40, h=0.92, w=12.25, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.10, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 color=DEEP, align=PP_ALIGN.LEFT, stroke_pt=1.5):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=stroke_pt,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.20)


def teal_callout(slide, x, y, w, h, text, *, size=14, bold=False,
                 color=DEEP, align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.18)


def footer(slide, text):
    text_box(slide, x=0.55, y=7.04, w=12.25, h=0.34, text=text,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             line_spacing=1.0)


def stat_plate(slide, x, y, w, h, label, value, *, highlight=False,
               vsize=25, lsize=12.5):
    bg = GOLD_TINT if highlight else SURFACE
    edge = GOLD if highlight else SOFT_GREY
    filled_rect(slide, x, y, w, h, bg, stroke=edge,
                stroke_pt=(1.5 if highlight else 1.0),
                radius=True, radius_adj=0.10)
    text_box(slide, x + 0.22, y + 0.08, w - 0.44, 0.30, label,
             size=lsize, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, x + 0.22, y + 0.34, w - 0.44, h - 0.42, value,
             size=vsize, bold=True, color=(GOLD if highlight else DEEP),
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


def icon(slide, name, x, y, size, variant="mid"):
    add_image(slide, ICONS / f"{name}-{variant}.png", x, y, size, size)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Deck loader — deck.yaml split на 2 части (≤600 строк каждая).
# Loader читает ОБЕ части, объединяет slides, вставляет s04a после s04.
# ============================================================
def load_deck():
    try:
        import yaml
    except ImportError:
        return None
    d1 = yaml.safe_load((ROOT / "deck.yaml").read_text(encoding="utf-8"))
    d2 = yaml.safe_load((ROOT / "deck-part2.yaml").read_text(encoding="utf-8"))
    slides = list(d1.get("slides", [])) + list(d2.get("slides", []))
    ids = [s["id"] for s in slides]
    expected = [f"s{n:02d}" for n in range(1, 33)]
    expected.insert(4, "s04a")  # s04a after s04
    assert ids == expected, (
        f"deck slide order mismatch:\n got={ids}\n exp={expected}")
    tot = d2.get("totals", {}).get("slides")
    assert tot == 33, f"deck-part2 totals.slides={tot}, expected 33"
    return {"slides": slides, "totals": d2.get("totals", {}),
            "deck": d1.get("deck", {})}


# ============================================================
# Section divider — unified template (7-card roadmap, gold current)
# Sections of Лекции 5: 0..6.
# ============================================================
NAV = [
    ("0", "Opening"),
    ("1", "Forecast"),
    ("2", "Anomalies"),
    ("3", "Scoring"),
    ("4", "LLM"),
    ("5", "Recomm."),
    ("6", "Framework"),
]


def roadmap_bar(slide, here_idx, *, y=6.50):
    """7-card progress bar; current section gold-bordered."""
    n = len(NAV)
    gap = 0.12
    bx = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    ch = 0.60
    for i, (num, label) in enumerate(NAV):
        x = bx + i * (cw + gap)
        cur = (i == here_idx)
        if cur:
            filled_rect(slide, x, y, cw, ch, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.12)
        else:
            filled_rect(slide, x, y, cw, ch, SURFACE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(slide, x=x + 0.04, y=y + 0.07, w=cw - 0.08, h=0.22,
                 text=f"Section {num}", size=10, bold=True,
                 color=(DEEP if cur else LIGHT), align=PP_ALIGN.CENTER)
        text_box(slide, x=x + 0.04, y=y + 0.30, w=cw - 0.08, h=0.26,
                 text=label, size=10.5, bold=cur,
                 color=(DEEP if cur else SLATE), align=PP_ALIGN.CENTER,
                 line_spacing=0.95)


def build_section_divider(p, here_idx, subtitle, bridge, sid):
    """Distinct divider (NO ocean motif): giant decorative section digit on
    the right, РАЗДЕЛ N + subtitle + 1-line narrative bridge on the left,
    gold-current roadmap bar at bottom."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=8.45, y=0.30, w=4.6, h=6.0, text=str(here_idx),
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.75, y=1.55, w=7.4, h=0.55,
             text=f"SECTION {here_idx}", size=20, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.18, 0.70, 0.05, fill=GOLD)
    text_box(s, x=0.75, y=2.55, w=7.6, h=1.85, text=subtitle,
             size=36, bold=True, color=DEEP, line_spacing=1.08)
    text_box(s, x=0.78, y=4.62, w=7.5, h=1.55, text=bridge,
             size=18, italic=True, color=LIGHT, line_spacing=1.22)
    roadmap_bar(s, here_idx, y=6.50)
    speaker_notes(s, load_notes(sid))
    return s


# ============================================================
# Slide builders — 33 slides (32 LOCKED s01..s32 + s04a divider)
# ============================================================

def build_s01(p):
    """case_study — Zillow iBuying hook. Left: 3 numbers in ocean box.
    Right: frame «какой тип ИИ и почему не LLM». Gold: $500M+ закрыто."""
    s = blank(p)
    slide_title(s, "A single predictive model shut down an entire business "
                   "line of a major company.", size=23, w=11.0, h=0.96)
    icon(s, "landmark", 12.05, 0.42, 0.74, "mid")
    # LEFT — mega-stat dominates; 3 small support stats below
    lx, ly, lw, lh = 0.55, 1.46, 7.05, 4.55
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.30, ly + 0.18, lw - 0.60, 0.30,
             "Zillow Offers · iBuying · November 2021",
             size=14, bold=True, color=MID)
    # mega number — the single point of impact
    filled_rect(s, lx + 0.30, ly + 0.54, lw - 0.60, 1.78, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.0, radius=True, radius_adj=0.08)
    text_box(s, lx + 0.50, ly + 0.66, lw - 1.00, 1.02, "$500M+",
             size=68, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.50, ly + 1.70, lw - 1.00, 0.52,
             "total losses — an entire business line shut down",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    sup = [("$300M+", "write-down for the quarter"),
           ("~2000", "laid off (≈25%)"),
           ("−25%", "stock in days")]
    n = 3
    gap = 0.18
    cw = (lw - 0.60 - gap * (n - 1)) / n
    sx = lx + 0.30
    for val, lab in sup:
        filled_rect(s, sx, ly + 2.50, cw, 0.94, SURFACE, stroke=SOFT_GREY,
                    stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, sx + 0.06, ly + 2.58, cw - 0.12, 0.42, val,
                 size=21, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, sx + 0.06, ly + 3.00, cw - 0.12, 0.38, lab,
                 size=11, color=MID, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        sx += cw + gap
    text_box(s, lx + 0.30, ly + 3.62, lw - 0.60, 0.78,
             "The model forecast the price of a house · Zillow "
             "automatically bought, renovated, and resold thousands of "
             "houses · systematically overvalued them",
             size=11.5, color=SLATE, line_spacing=1.16)
    # RIGHT — "не ChatGPT" + open question
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 2.30, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "circle-help", rx + 0.28, ly + 0.24, 0.50, "teal")
    text_box(s, rx + 0.92, ly + 0.28, rw - 1.10, 0.42, "This was not ChatGPT",
             size=17, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.28, ly + 0.90, rw - 0.56, 1.28,
             "A predictive model estimated a number (a price) from tabular "
             "and geo-data. An LLM was not — and could not be — used here: "
             "the task is not textual.",
             size=13.5, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, ly + 2.46, rw, 2.09)
    text_box(s, rx + 0.26, ly + 2.62, rw - 0.52, 1.80,
             "What type of AI was this — and why did an ordinary model error "
             "become a business loss, rather than a minor inaccuracy?",
             size=14.5, italic=True, color=MID, line_spacing=1.22,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 6.20, 12.25, 0.66,
                 "A forecast error is normal. The business collapsed because "
                 "of WHERE it was wired: to an automated, irreversible action "
                 "at the scale of thousands of houses.", size=14)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """cover — distinct, NO ocean motif. Mega «05» + title + roadmap-bar."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=7.7, y=0.95, w=5.7, h=5.2, text="05",
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.75, y=1.45, w=6.6, h=0.5, text="LECTURE 5",
             size=18, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.02, 0.70, 0.05, fill=TEAL)
    text_box(s, x=0.75, y=2.42, w=7.7, h=2.5,
             text="AI in Finance\nand Retail",
             size=44, bold=True, color=DEEP, line_spacing=1.10)
    filled_rect(s, 0.78, 5.18, 0.05, 0.56, fill=TEAL)
    text_box(s, x=1.02, y=5.16, w=7.6, h=0.62,
             text="For which task — which type of AI,\nwhy that one, and "
                  "where it breaks",
             size=18, color=MID, line_spacing=1.18)
    roadmap_bar(s, 0, y=6.50)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    """comparison — bridge Л2/Л4 → палитра 5 типов. KEYSTONE."""
    s = blank(p)
    slide_title(s, "One type of AI in depth (L4) → a palette of types for a "
                   "palette of tasks (L5).", size=24, y=0.32, h=0.56, w=12.25)
    # compact secondary bridge strip (Л2/Л4 → Л5) — single line each
    bx, by, bw, bh = 0.55, 0.98, 12.25, 1.16
    ocean_box(s, bx, by, bw, bh)
    bridges = [
        ("L2: the «which AI» tree — ML / deep learning / LLM",
         "L5: five structurally different types for different tasks"),
        ("L4: one type (LLM coder) in depth along autonomy",
         "L5: a palette of types for a palette of tasks in one industry"),
    ]
    yy = by + 0.16
    for a, b in bridges:
        text_box(s, bx + 0.26, yy, 5.55, 0.40, a,
                 size=11.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        text_box(s, bx + 5.86, yy + 0.06, 0.34, 0.30, "→",
                 size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, bx + 6.26, yy, bw - 6.52, 0.40, b,
                 size=11.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        yy += 0.44
    # MAIN VISUAL — palette of 5 types, large
    text_box(s, 0.55, 2.34, 12.25, 0.34,
             "Map of five AI types for five different tasks (+ a "
             "cross-cutting CV layer):", size=15, bold=True, color=MID)
    types = ["Time-series\nforecasting", "Anomaly\ndetection", "ML\nscoring",
             "LLM\nassistants", "Recommendations\nand pricing"]
    icns = ["trending-up", "radar", "scale", "message-circle",
            "shopping-cart"]
    n = 5
    gap = 0.22
    cw = (12.25 - gap * (n - 1)) / n
    ax, ay, ah = 0.55, 2.78, 2.06
    for i in range(n):
        ocean_box(s, ax, ay, cw, ah)
        icon(s, icns[i], ax + cw / 2 - 0.36, ay + 0.30, 0.72, "mid")
        for li, line in enumerate(types[i].split("\n")):
            text_box(s, ax + 0.04, ay + 1.18 + li * 0.36, cw - 0.08, 0.36,
                     line, size=14, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.0)
        ax += cw + gap
    gold_callout(s, 0.55, 5.10, 12.25, 1.66,
                 "Most of the value here comes NOT from a language model. "
                 "The LLM is not a universal hammer: a different task "
                 "requires a different type of AI. That is exactly «a "
                 "palette of types for a palette of tasks».", size=16)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """assertion_visual — central question + 5-step pattern + disclaimer."""
    s = blank(p)
    slide_title(s, "The central question of the lecture.",
                size=27, y=0.34, h=0.58, w=9.5)
    icon(s, "target", 12.05, 0.36, 0.78, "gold")
    bx, by, bw, bh = 0.55, 1.18, 12.25, 1.62
    ocean_box(s, bx, by, bw, bh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, bx + 0.40, by + 0.18, bw - 0.80, bh - 0.36,
             "«Finance and retail are the industries with the highest AI "
             "adoption. For which task — which type of AI, why that one (and "
             "not LLM everywhere), and where does this type break?»",
             size=20, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16, align=PP_ALIGN.CENTER)
    # 5-step pattern ribbon
    text_box(s, 0.55, 2.96, 12.25, 0.30,
             "We examine each type by one scheme (5 steps):",
             size=14, bold=True, color=MID)
    steps = ["task", "type of AI\nand why it", "real\nexample",
             "where it\nbreaks", "alternative\nand criterion"]
    n = 5
    gap = 0.18
    cw = (12.25 - gap * (n - 1)) / n
    sx, sy, sh = 0.55, 3.30, 0.92
    for i in range(n):
        filled_rect(s, sx, sy, cw, sh, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                    radius=True, radius_adj=0.12)
        text_box(s, sx + 0.06, sy, cw - 0.12, sh, f"{i + 1}. {steps[i]}",
                 size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        if i < n - 1:
            text_box(s, sx + cw - 0.02, sy, 0.20, sh, "→",
                     size=16, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        sx += cw + gap
    teal_callout(s, 0.55, 4.42, 12.25, 0.66,
                 "The map is the structure of the lecture, not a requirement "
                 "to understand everything now. Forecasting · anomalies · "
                 "scoring · LLM · recommendations (+ a cross-cutting CV).",
                 size=13.5, bold=True)
    gold_callout(s, 0.55, 5.26, 12.25, 1.50,
                 "The answer is not «AI is good» and not «AI is bad», but an "
                 "apparatus: name the appropriate type, justify the choice, "
                 "and see in advance the point where a human is mandatory. "
                 "By the end this scheme will become your tool of choice.",
                 size=14)
    speaker_notes(s, load_notes("s04"))


def build_s04a(p):
    build_section_divider(
        p, 1, "Forecasting:\ndemand, sales, churn",
        "We start with retail's most widespread task — and right away with a "
        "type of AI that is NOT an LLM.", "s04a")


def build_s05(p):
    """assertion_visual — prognosis vs LLM, 2-col parallel."""
    s = blank(p)
    slide_title(s, "Time-series forecasting and a language model solve "
                   "structurally different tasks.", size=24, y=0.34, h=0.62,
                w=12.25)
    bx, by, bw, bh = 0.55, 1.20, 12.25, 3.40
    ocean_box(s, bx, by, bw, bh)
    col_w = [6.0, 6.05]
    heads = ["Language model (LLM)", "Time-series forecasting"]
    cx = bx + 0.20
    for j, hd in enumerate(heads):
        filled_rect(s, cx, by + 0.18, col_w[j] - 0.20, 0.48,
                    MID if j == 0 else TEAL)
        text_box(s, cx, by + 0.18, col_w[j] - 0.20, 0.48, hd,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    rows = [
        ("predicts the next TOKEN of text",
         "predicts the next NUMBER in the series"),
        ("optimized for the plausibility of the wording",
         "needs numeric accuracy + calibrated uncertainty"),
        ("has no internal notion of «seasonality» / «trend»",
         "series = trend + seasonality + noise"),
        ("has seen texts ABOUT sales",
         "sees the sales series itself as a series"),
    ]
    yy = by + 0.78
    for a, b in rows:
        cx = bx + 0.20
        for j, cc in enumerate((a, b)):
            filled_rect(s, cx, yy, col_w[j] - 0.20, 0.58,
                        SURFACE, stroke=SOFT_GREY, stroke_pt=0.75)
            text_box(s, cx + 0.16, yy, col_w[j] - 0.42, 0.58, cc,
                     size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.05)
            cx += col_w[j]
        yy += 0.62
    teal_callout(s, 0.55, 4.74, 12.25, 0.62,
                 "Type of AI for forecasting: classical statistics and "
                 "tabular ML — the ARIMA family, gradient boosting. NOT a "
                 "generative model, NOT an LLM.",
                 size=13.5, bold=True)
    gold_callout(s, 0.55, 5.52, 12.25, 1.26,
                 "Hand «forecast the demand» to a text model — it will "
                 "produce a plausible number with no connection to "
                 "seasonality: a hallucination in a numeric wrapper. The "
                 "type of AI is determined by the STRUCTURE of the task, not "
                 "by the trendiness of the tool.", size=14)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """assertion_visual — 3 argument cards (why not LLM) + alternative."""
    s = blank(p)
    slide_title(s, "Why a tabular model here, not an LLM — three arguments, "
                   "not a caveat.", size=22, y=0.34, h=0.58,
                w=12.25)
    cards = [
        ("database", "1. Data structure",
         "tabular numeric series (date, point, product, price, features). "
         "Series models extract trend and seasonality; turning a series "
         "into text for an LLM loses exactly this structure.", "mid", False),
        ("activity", "2. Measurability and calibration",
         "the decision about order size is built on a numeric error metric "
         "+ a confidence interval. Classical models give this out of the "
         "box; an LLM does not.", "mid", False),
        ("triangle-alert", "3. What will break",
         "an LLM will give a plausible number with no grounded uncertainty "
         "— and the decision still has to be made, and the cost of a "
         "systematic error × volume (the Zillow mechanism).", "gold", True),
    ]
    n = 3
    gap = 0.26
    cw = (12.25 - gap * (n - 1)) / n
    cx = 0.55
    cy, chh = 1.22, 2.74
    for ic, ttl, body, var, hi in cards:
        if hi:
            ocean_box(s, cx, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD,
                      stroke_pt=2.0)
        else:
            ocean_box(s, cx, cy, cw, chh)
        icon(s, ic, cx + 0.24, cy + 0.22, 0.50, var)
        text_box(s, cx + 0.86, cy + 0.26, cw - 1.05, 0.50, ttl,
                 size=15, bold=True, color=(DEEP if hi else MID),
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + 0.26, cy + 0.90, cw - 0.52, chh - 1.02, body,
                 size=12, color=DEEP, line_spacing=1.16)
        cx += cw + gap
    teal_callout(s, 0.55, 4.12, 12.25, 0.84,
                 "Alternative and criterion: the right tool is a tabular "
                 "predictive model with a measurable error. Even it is not "
                 "enough under a non-stationary environment + an irreversible "
                 "capital action without human control.",
                 size=13.5, bold=True)
    gold_callout(s, 0.55, 5.08, 12.25, 0.86,
                 "Tasks: demand forecasting (shelf / procurement), cash-flow "
                 "(liquidity), customer churn (whom to retain).", size=14)
    footer(s, "If the series contain PII of Russian citizens, sending them to "
              "a public cloud LLM violates localization requirements "
              "(Federal Law 152-FZ); a tabular model is deployed on-premise.")
    speaker_notes(s, load_notes("s06"))


def build_s07(p):
    """case_study — X5 + Магнит, 2 panels."""
    s = blank(p)
    slide_title(s, "The forecast is dictated by the task: neither X5 nor "
                   "Magnit chose an LLM.", size=24, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.30, 6.10, 3.95
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.34,
             "X5 Group · «Pyaterochka», «Perekrestok»",
             size=14, bold=True, color=MID)
    x5 = [("Demand forecast accuracy", "> 70%", False),
          ("Extra revenue from ML tools", "+5 bn ₽", True),
          ("Expiry write-offs", "−2%", False)]
    cyy = ly + 0.58
    for lab, val, hi in x5:
        bg = GOLD_TINT if hi else SURFACE
        filled_rect(s, lx + 0.26, cyy, lw - 0.52, 0.96, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.44, cyy + 0.10, lw - 2.30, 0.76, lab,
                 size=12.5, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
        text_box(s, lx + lw - 2.05, cyy + 0.06, 1.66, 0.84, val,
                 size=22, bold=True, color=(GOLD if hi else DEEP),
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        cyy += 1.04
    rx, rw = 6.85, 5.95
    ocean_box(s, rx, ly, rw, 3.95)
    text_box(s, rx + 0.26, ly + 0.16, rw - 0.52, 0.34,
             "Magnit — import substitution of forecasting systems",
             size=14, bold=True, color=MID)
    mg = [
        ("Before 2022:", "foreign vendors of the SAP / Blue Yonder class"),
        ("After they left:", "own forecasting + auto-ordering system (F&R)"),
        ("Status:", "pilot at a distribution center (2024–2025)"),
    ]
    fy = ly + 0.62
    for a, b in mg:
        text_box(s, rx + 0.26, fy, 1.55, 0.66, a, size=12.5, bold=True,
                 color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 1.86, fy, rw - 2.12, 0.66, b, size=12.5, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        fy += 0.74
    teal_callout(s, rx, ly + 2.92, rw, 0.90,
                 "The departure of foreign vendors made import substitution "
                 "of forecasting systems a direct engineering task for the "
                 "industry.",
                 size=12.5, bold=False)
    gold_callout(s, 0.55, 5.40, 12.25, 0.88,
                 "Neither X5 nor Magnit solves the task with a language model "
                 "— they build specialized forecasting systems, because the "
                 "TYPE OF AI IS DICTATED BY THE TASK.", size=14)
    footer(s, "Per company data and industry reviews (X5/TAdviser 2023; "
              "shoppers.media/TAdviser 2024–2025); the numbers are as claimed "
              "by the companies, verified on the day of the lecture.")
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """process — time-series decomposition diagram + товаровед analogy."""
    s = blank(p)
    slide_title(s, "A forecast extends the patterns of the past — here lie "
                   "both its strength and its vulnerability.", size=23,
                y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.22, 7.30, 3.95
    ocean_box(s, lx, ly, lw, lh)
    add_image(s, DIAGRAMS / "d08-timeseries-decomp.png",
              lx + 0.20, ly + 0.16, lw - 0.40, lh - 0.62)
    text_box(s, lx + 0.24, ly + lh - 0.42, lw - 0.48, 0.34,
             "Any sales series = trend + seasonality + noise; the model "
             "extends the regular part.",
             size=10.5, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    rx, rw = 8.05, 4.75
    ocean_box(s, rx, ly, rw, 3.95)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.30,
             "What the model does — numerically", size=13.5, bold=True,
             color=MID)
    pts = [
        "decomposes history into trend + seasonality + noise",
        "learns the regular patterns and extends them",
        "adds corrections for promo / holiday",
    ]
    py = ly + 0.54
    for t in pts:
        circle(s, rx + 0.26, py + 0.05, 0.12, MID)
        text_box(s, rx + 0.52, py, rw - 0.78, 0.56, t,
                 size=12, color=DEEP, line_spacing=1.10)
        py += 0.58
    text_box(s, rx + 0.24, py + 0.06, rw - 0.48, 0.46,
             "Forecast → action: «demand ≈ N» → «order N + buffer»",
             size=12, bold=True, color=TEAL, line_spacing=1.08)
    text_box(s, rx + 0.24, py + 0.62, rw - 0.48, 1.00,
             "Analogy: a forecasting model is like an experienced buyer "
             "«stock up more for the weekend», but for millions of "
             "store × product pairs at once.",
             size=11.5, italic=True, color=SLATE, line_spacing=1.16)
    gold_callout(s, 0.55, 5.32, 12.25, 0.96,
                 "Built-in vulnerability: a forecast is strong exactly as "
                 "long as TOMORROW RESEMBLES YESTERDAY. When the environment "
                 "changes qualitatively — the model confidently extrapolates "
                 "patterns that no longer exist. The key to the next slide.",
                 size=14)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """case_study — Zillow payoff: distribution shift + asymmetry + Knight."""
    s = blank(p)
    slide_title(s, "What ruined Zillow was not model drift, but what its "
                   "output was wired to.", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.16, 7.05, 4.10
    ocean_box(s, lx, ly, lw, lh)
    blocks = [
        ("distribution shift",
         "the real data stopped resembling the training data — the housing "
         "market of 2020–2021 became statistically different"),
        ("asymmetry of the cost of error",
         "a recommendation error ≈ 0 · the same error on which a house is "
         "bought = tens of thousands of $ × N, and IRREVERSIBLE"),
        ("Three things together = fatal",
         "irreversible action × automated on the model's output × a "
         "non-stationary environment without a circuit-breaker (auto-stop "
         "on anomaly)"),
    ]
    by = ly + 0.20
    for ttl, body in blocks:
        text_box(s, lx + 0.26, by, lw - 0.52, 0.28, ttl,
                 size=13, bold=True, color=MID)
        text_box(s, lx + 0.26, by + 0.30, lw - 0.52, 0.60, body,
                 size=12, color=DEEP, line_spacing=1.12)
        by += 1.02
    filled_rect(s, lx + 0.26, by + 0.04, lw - 0.52, 0.62, TEAL_TINT,
                stroke=TEAL, stroke_pt=1.2, radius=True, radius_adj=0.14)
    text_box(s, lx + 0.42, by + 0.04, lw - 0.84, 0.62,
             "Knight Capital, 2012: $440M / ~45 min on a deterministic "
             "algorithm — the same class of error",
             size=11.5, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 1.62)
    text_box(s, rx + 0.26, ly + 0.14, rw - 0.52, 0.30,
             "Zillow vs Opendoor", size=14, bold=True, color=MID)
    text_box(s, rx + 0.26, ly + 0.48, rw - 0.52, 1.02,
             "The same type of AI, the same period. Opendoor survived "
             "thanks to a more conservative spread and risk design.",
             size=12.5, color=DEEP, line_spacing=1.16)
    teal_callout(s, rx, ly + 1.78, rw, 2.32,
                 "It is dangerous when simultaneously: (a) the output "
                 "triggers a large irreversible action automatically, (b) the "
                 "environment is non-stationary, (c) there is no "
                 "circuit-breaker.\n\nAlternative: a narrow segment + "
                 "human-gate + live error-monitoring + circuit-breaker.",
                 size=12, bold=False)
    gold_callout(s, 0.55, 5.40, 12.25, 0.90,
                 "The model did exactly what it was meant for. The error was "
                 "the engineering decision of WHERE to wire its output. The "
                 "same AI, a different judgment → bankruptcy vs survival.",
                 size=14)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    build_section_divider(
        p, 2, "Anomaly detection:\nfraud and AML",
        "Forecasting is about the future; now about the present — catch an "
        "anomaly within milliseconds, while the payment has not yet gone "
        "through.", "s10")


def build_s11(p):
    """assertion_visual — anomaly cloud diagram + why-not."""
    s = blank(p)
    slide_title(s, "Learn the «customer's norm» and catch the deviation — "
                   "not learn «what fraud looks like».", size=22, y=0.34,
                h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.20, 5.55, 4.05
    ocean_box(s, lx, ly, lw, lh)
    add_image(s, DIAGRAMS / "d11-anomaly-cloud.png",
              lx + 0.20, ly + 0.18, lw - 0.40, lh - 0.36)
    rx, rw = 6.30, 6.50
    ocean_box(s, rx, ly, rw, 4.05)
    text_box(s, rx + 0.26, ly + 0.16, rw - 0.52, 0.62,
             "Task: catch a fraudulent transaction in real time, where "
             "almost everything is legitimate and the share of fraud is "
             "extremely small.",
             size=13, bold=True, color=MID, line_spacing=1.14)
    text_box(s, rx + 0.26, ly + 0.86, rw - 0.52, 0.56,
             "Type of AI — anomaly detection: the model builds the "
             "customer's «norm» and signals deviations.",
             size=12.5, color=DEEP, line_spacing=1.14)
    whys = [
        ("not classification", "fraud is rare and constantly changes its "
                               "form (a moving target)"),
        ("not time-series forecasting", "the question is not «the next "
                            "value», but «how unlike the norm»"),
        ("not an LLM", "the transaction is structured; the task is "
                   "geometric, not linguistic"),
    ]
    fy = ly + 1.52
    for a, b in whys:
        text_box(s, rx + 0.26, fy, 1.85, 0.56, a, size=12.5, bold=True,
                 color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 2.16, fy, rw - 2.42, 0.56, b, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        fy += 0.60
    filled_rect(s, rx + 0.26, fy + 0.06, rw - 0.52, 0.62, SURFACE,
                stroke=SOFT_GREY, stroke_pt=0.75, radius=True,
                radius_adj=0.12)
    text_box(s, rx + 0.40, fy + 0.06, rw - 0.80, 0.62,
             "AML is a subset of the same task; hard statutory thresholds = "
             "deterministic rules, not a model.",
             size=11, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "The model does not «know» that this is fraud — it knows "
                 "that this is NOT LIKE THE NORM, and raises a flag for "
                 "review.", size=14)
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """case_study — fraud chart + hint."""
    s = blank(p)
    slide_title(s, "Anti-fraud is «norm / deviation» in real time, not "
                   "generation.", size=23, y=0.34, h=0.58, w=12.25)
    # 3 SEPARATE stat-plates — each its own unit (no false shared axis)
    text_box(s, 0.55, 1.16, 12.25, 0.30,
             "What anomaly detection delivers in anti-fraud — each metric in "
             "its own units:", size=13.5, bold=True, color=MID)
    plates = [
        ("Stripe Radar", "−32% fraud", "while approving legitimate > 99%",
         False),
        ("JPMorgan", "−30% false", "fewer false positives",
         False),
        ("Visa · prevented", "~$40 bn", "of fraudulent transactions "
         "(FY2023)", True),
    ]
    n = 3
    gap = 0.24
    cw = (12.25 - gap * (n - 1)) / n
    px, py, ph = 0.55, 1.54, 2.16
    for lab, val, sub, hi in plates:
        bg = GOLD_TINT if hi else SURFACE
        edge = GOLD if hi else LIGHT
        ocean_box(s, px, py, cw, ph, fill=bg, stroke=edge,
                  stroke_pt=(2.0 if hi else 1.5))
        text_box(s, px + 0.20, py + 0.20, cw - 0.40, 0.34, lab,
                 size=14, bold=True, color=MID, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, px + 0.16, py + 0.62, cw - 0.32, 0.78, val,
                 size=33, bold=True, color=(GOLD if hi else DEEP),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, px + 0.20, py + 1.46, cw - 0.40, 0.58, sub,
                 size=12, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        px += cw + gap
    # Russia context + the metric hint, side by side below plates
    lx2, lw2 = 0.55, 6.00
    ocean_box(s, lx2, 3.92, lw2, 1.04)
    text_box(s, lx2 + 0.24, 4.04, lw2 - 0.48, 0.28, "Russia",
             size=13, bold=True, color=MID)
    text_box(s, lx2 + 0.24, 4.34, lw2 - 0.48, 0.56,
             "Anomaly detection is standard practice at major banks "
             "(per Bank of Russia materials, 20.11.2025).",
             size=11.5, color=DEEP, line_spacing=1.12)
    rx2, rw2 = 6.80, 6.00
    teal_callout(s, rx2, 3.92, rw2, 1.04,
                 "Notice: «reduction of FALSE positives» — the key metric of "
                 "anti-fraud is NOT «accuracy in general», but the ratio of "
                 "the two types of error. We will examine this further.",
                 size=11.5, bold=False)
    gold_callout(s, 0.55, 5.12, 12.25, 0.86,
                 "Different units — NOT one scale: % reduction, % false, "
                 "bn $. And still the main metric is not «accuracy», but "
                 "FP / FN separately.", size=14)
    footer(s, "Stripe/JPMorgan/Visa — as claimed by the companies; Visa — per "
              "Reuters/CNBC reports, July 2024; Russia — per Bank of Russia "
              "materials, 20.11.2025; the numbers are verified on the day of "
              "the lecture.")
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """comparison / schema_matrix — confusion matrix 2x2 + accuracy lie."""
    s = blank(p)
    slide_title(s, "«99.9% accuracy» under strong imbalance is deceptive.",
                size=24, y=0.34, h=0.58, w=12.25)
    bx, by, bw, bh = 0.55, 1.18, 7.55, 3.10
    ocean_box(s, bx, by, bw, bh)
    # 2x2 matrix
    mx, my = bx + 0.30, by + 0.24
    col_lab = ["Was fraud", "Was legitimate"]
    cw2, ch2 = 3.30, 1.10
    text_box(s, mx + 1.10, my - 0.02, cw2, 0.26, col_lab[0],
             size=11.5, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(s, mx + 1.10 + cw2, my - 0.02, cw2, 0.26, col_lab[1],
             size=11.5, bold=True, color=MID, align=PP_ALIGN.CENTER)
    rlab = ["System:\nfraud", "System:\npass"]
    cells = [
        [("TP", "caught (good)", TEAL_TINT, TEAL),
         ("FP", "honest customer blocked\n(type I error)",
          GOLD_TINT, GOLD)],
        [("FN", "money went to the fraudster\n(type II error)",
          RGBColor(0xE6, 0xE9, 0xF2), DEEP),
         ("TN", "correctly passed\n(good)", TEAL_TINT, TEAL)],
    ]
    for ri in range(2):
        text_box(s, mx - 0.04, my + 0.30 + ri * (ch2 + 0.12), 1.05, ch2,
                 rlab[ri], size=11, bold=True, color=MID,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                 line_spacing=1.0)
        for ci in range(2):
            tag, desc, fill, fg = cells[ri][ci]
            xx = mx + 1.10 + ci * (cw2 + 0.05)
            yy = my + 0.30 + ri * (ch2 + 0.12)
            filled_rect(s, xx, yy, cw2, ch2, fill,
                        stroke=fg, stroke_pt=1.2, radius=True,
                        radius_adj=0.08)
            text_box(s, xx + 0.12, yy + 0.08, cw2 - 0.22, 0.32, tag,
                     size=16, bold=True, color=fg)
            text_box(s, xx + 0.12, yy + 0.44, cw2 - 0.22, ch2 - 0.52, desc,
                     size=11, color=DEEP, line_spacing=1.05)
    rx, rw = 8.30, 4.50
    ocean_box(s, rx, by, rw, 3.10)
    text_box(s, rx + 0.24, by + 0.16, rw - 0.48, 0.30,
             "Why accuracy lies", size=13.5, bold=True, color=MID)
    text_box(s, rx + 0.24, by + 0.52, rw - 0.48, 1.10,
             "A stream of 1,000,000, fraud 1000 (0.1%). A model that says "
             "«everything is legitimate» gives 99.9% accuracy — and is "
             "completely useless.",
             size=12, color=DEEP, line_spacing=1.18)
    filled_rect(s, rx + 0.24, by + 1.66, rw - 0.48, 0.02, SOFT_GREY)
    text_box(s, rx + 0.24, by + 1.78, rw - 0.48, 0.30,
             "cost-sensitive", size=13.5, bold=True, color=MID)
    text_box(s, rx + 0.24, by + 2.12, rw - 0.48, 0.90,
             "FP and FN differ in COST (money, trust). Minimize the total "
             "expected cost, not a «pretty» accuracy.",
             size=12, color=DEEP, line_spacing=1.18)
    text_box(s, 0.55, 4.40, 12.25, 0.26,
             "The formal apparatus (sensitivity / specificity) is built in "
             "Lecture 7 on a medical example; here — the working intuition.",
             size=11, italic=True, color=SLATE)
    gold_callout(s, 0.55, 4.78, 12.25, 0.80,
                 "The first question to any anti-fraud number is NOT «what is "
                 "the accuracy», but «what are the FP and FN separately, and "
                 "at what cost».",
                 size=14)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """case_study — FP cost contrast + precision/recall + Knight 2."""
    s = blank(p)
    slide_title(s, "«We cut FP by 25%» averages a penny FP and a "
                   "catastrophic FP.", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 7.05, 4.10
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.30,
             "The same class of error — cost differs by orders of magnitude",
             size=13.5, bold=True, color=MID)
    fp = [
        ("FP #1 · $5 for coffee", "the customer retried, forgot within an "
                               "hour → cost ≈ 0, REVERSIBLE", False),
        ("FP #2 · $5000 for treatment abroad",
         "hard to get through, time is critical → IRREVERSIBLE in its "
         "consequences",
         True),
    ]
    fy = ly + 0.52
    for ttl, body, hi in fp:
        bg = GOLD_TINT if hi else SURFACE
        filled_rect(s, lx + 0.26, fy, lw - 0.52, 0.86, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.42, fy + 0.08, lw - 0.84, 0.30, ttl,
                 size=12.5, bold=True, color=(DEEP if hi else MID))
        text_box(s, lx + 0.42, fy + 0.40, lw - 0.84, 0.42, body,
                 size=11, color=DEEP, line_spacing=1.05)
        fy += 0.94
    text_box(s, lx + 0.26, fy + 0.04, lw - 0.52, 0.74,
             "precision ↔ recall — a TRADE-OFF, not a task of «both to "
             "zero». A better model shifts the curve; the point is chosen by "
             "the engineer per the cost of error.",
             size=11.5, bold=True, color=TEAL, line_spacing=1.12)
    filled_rect(s, lx + 0.26, fy + 0.84, lw - 0.52, 0.46, TEAL_TINT,
                stroke=TEAL, stroke_pt=1.2, radius=True, radius_adj=0.16)
    text_box(s, lx + 0.40, fy + 0.84, lw - 0.80, 0.46,
             "Auto-block of the irreversible without a gate — the Zillow / "
             "Knight class",
             size=11, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    rx, rw = 7.80, 5.00
    teal_callout(s, rx, ly, rw, 4.10,
                 "Criterion:\n\n• auto-hard block — only for the reversible / "
                 "small\n\n• large / irreversible → a soft challenge "
                 "(3DS, call, push) + a fast human channel for "
                 "unblocking\n\n• a hard AML threshold — a rules-engine "
                 "(law, not «with probability 0.97»)\n\nThe cost of FP "
                 "averages the penny and the catastrophic — measure them "
                 "separately.",
                 size=13, bold=False)
    gold_callout(s, 0.55, 5.40, 12.25, 0.90,
                 "Scale does not reduce harm: 0.5% false of billions = tens "
                 "of millions of blocked legitimate transactions. Behind "
                 "each one — a specific person.", size=14)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    build_section_divider(
        p, 3, "Credit scoring:\nwhy NOT a neural network",
        "The anomaly has been found; now a decision that changes the "
        "customer's life — and why a black box is not allowed here.", "s15")


def build_s16(p):
    """assertion_visual — 4 argument cards + инспектор analogy."""
    s = blank(p)
    slide_title(s, "A black box in scoring is inapplicable in principle — not "
                   "«less convenient».", size=23, y=0.34, h=0.58, w=12.25)
    cards = [
        ("gavel", "1. Explainability = law",
         "a customer who is declined is legally owed an explanation (reason "
         "codes: «high debt burden»). «The algorithm decided so» is "
         "impermissible.",
         "mid", False),
        ("database", "2. The data is tabular",
         "deep networks are strong on text/audio/image; on tabular data, "
         "boosting is competitive at incomparably greater interpretability.",
         "mid", False),
        ("badge-check", "3. Audit > +1% accuracy",
         "the regulator must reproduce the decision, verify the absence of "
         "discrimination. A stable, explainable one matters more than an "
         "opaque, more accurate one.",
         "mid", False),
        ("triangle-alert", "4. What breaks with a black box",
         "cannot be explained → a violation · cannot be audited for bias → "
         "a crisis (Apple Card) · unstable → cannot be defended before the "
         "regulator.", "gold", True),
    ]
    # left: 4 argument cards (compact, single column band) ; right: analogy
    cw2, ch2 = 5.95, 0.96
    for i, (ic, ttl, body, var, hi) in enumerate(cards):
        cx = 0.55
        cy = 1.16 + i * (ch2 + 0.14)
        if hi:
            ocean_box(s, cx, cy, cw2, ch2, fill=GOLD_TINT, stroke=GOLD,
                      stroke_pt=2.0)
        else:
            ocean_box(s, cx, cy, cw2, ch2)
        icon(s, ic, cx + 0.18, cy + (ch2 - 0.36) / 2, 0.36, var)
        text_box(s, cx + 0.66, cy + 0.10, cw2 - 0.84, 0.30, ttl,
                 size=12.5, bold=True, color=(DEEP if hi else MID))
        text_box(s, cx + 0.66, cy + 0.40, cw2 - 0.84, ch2 - 0.48, body,
                 size=10.5, color=DEEP, line_spacing=1.08)
    # right — inspector analogy diagram (chapter §3.3, derived)
    rx, ry, rw, rh = 6.70, 1.16, 6.10, 4.54
    ocean_box(s, rx, ry, rw, rh)
    text_box(s, rx + 0.24, ry + 0.14, rw - 0.48, 0.30,
             "Anchor analogy: a credit inspector", size=13, bold=True,
             color=MID)
    add_image(s, DIAGRAMS / "d16-inspector-reason-codes.png",
              rx + 0.18, ry + 0.50, rw - 0.36, rh - 1.30)
    text_box(s, rx + 0.24, ry + rh - 0.74, rw - 0.48, 0.62,
             "An interpretable model shows the calculation line by line "
             "(reason codes); a black box refuses to explain.",
             size=10.5, italic=True, color=SLATE, line_spacing=1.10)
    gold_callout(s, 0.55, 5.86, 12.25, 0.92,
                 "«New = a neural network, so better» in scoring is a "
                 "structural error: explainability here is not a bonus, but a "
                 "CONDITION OF LEGALITY.", size=14)
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """case_study — Сбер stats + clarification."""
    s = blank(p)
    slide_title(s, "«100% AI decisions» in Russia ≠ a neural-network black "
                   "box out of control.", size=23, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.30, 5.55, 3.95
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.32,
             "Sberbank · since February–March 2024", size=13.5, bold=True,
             color=MID)
    sb = [("Decisions on individuals", "~100% AI", False),
          ("Scoring model", "up to 5000 parameters", False),
          ("AI effect across all lines, 2023", "+350 bn ₽", True)]
    cyy = ly + 0.56
    for lab, val, hi in sb:
        bg = GOLD_TINT if hi else SURFACE
        filled_rect(s, lx + 0.26, cyy, lw - 0.52, 1.00, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.42, cyy + 0.10, lw - 0.84, 0.42, lab,
                 size=12, bold=True, color=MID, line_spacing=1.05)
        text_box(s, lx + 0.42, cyy + 0.50, lw - 0.84, 0.44, val,
                 size=20, bold=True, color=(GOLD if hi else DEEP))
        cyy += 1.08
    rx, rw = 6.30, 6.50
    ocean_box(s, rx, ly, rw, 3.95)
    text_box(s, rx + 0.26, ly + 0.14, rw - 0.52, 0.28,
             "«100% AI» actually means:", size=14, bold=True,
             color=MID)
    items = [
        ("high automation of the pipeline on interpretable models",
         False),
        ("reason codes on every decision", False),
        ("regulatory supervision by the Bank of Russia", False),
        ("the customer's retained right to a human: > 80% of financial "
         "organizations offer a human opt-out", True),
    ]
    iy = ly + 0.48
    for it, hi in items:
        ih = 0.74 if hi else 0.52
        if hi:
            filled_rect(s, rx + 0.20, iy, rw - 0.40, ih, GOLD_TINT,
                        stroke=GOLD, stroke_pt=1.5, radius=True,
                        radius_adj=0.12)
            circle(s, rx + 0.38, iy + ih / 2 - 0.06, 0.13, GOLD)
            text_box(s, rx + 0.66, iy, rw - 0.92, ih, it,
                     size=12, bold=True, color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
        else:
            circle(s, rx + 0.30, iy + 0.13, 0.12, MID)
            text_box(s, rx + 0.56, iy, rw - 0.82, ih, it,
                     size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.10)
        iy += ih + 0.14
    teal_callout(s, rx, ly + 3.24, rw, 0.60,
                 "The type of AI — tabular ML, chosen because of "
                 "explainability, not in spite of it. The lecture's thesis at "
                 "the scale of the largest bank.",
                 size=11.5, bold=False)
    footer(s, "Per the bank's statements (TAdviser/AdIndex/Interfax, 2024) and "
              "Bank of Russia materials (20.11.2025); the numbers are verified "
              "on the day of the lecture.")
    speaker_notes(s, load_notes("s17"))


def build_s18(p):
    """case_study — Apple Card precise outcome + proxy-bias from scratch."""
    s = blank(p)
    slide_title(s, "Even a formally lawful model without explainability "
                   "creates a crisis.", size=23, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 6.35, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.32,
             "Apple Card / Goldman Sachs · 2019–2021", size=13.5,
             bold=True, color=MID)
    timeline = [
        "Nov 2019: a viral thread — the limit is ×20 higher for the husband "
        "(while the wife's credit rating is higher)",
        "NYDFS opened an investigation (~400,000 applicants)",
        "March 2021: NYDFS found NO violation of the law",
        "BUT: customers received no explanations → a regulatory-reputational "
        "crisis",
    ]
    ty = ly + 0.54
    for i, t in enumerate(timeline):
        hi = (i == 2)
        circle(s, lx + 0.26, ty + 0.06, 0.13, GOLD if hi else MID)
        text_box(s, lx + 0.52, ty, lw - 0.78, 0.74, t,
                 size=11.5, bold=hi, color=(DEEP if not hi else DEEP),
                 line_spacing=1.10)
        ty += 0.80
    text_box(s, lx + 0.26, ty + 0.02, lw - 0.52, 0.40,
             "To claim «discrimination was proven» is factually incorrect.",
             size=11, italic=True, color=SLATE,
             line_spacing=1.05)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 2.32)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.30,
             "Proxy bias «in plain terms»", size=13.5, bold=True,
             color=MID)
    text_box(s, rx + 0.24, ly + 0.48, rw - 0.48, 1.72,
             "The engineer does not feed in gender/race. But features remain "
             "that CORRELATE with them: postal code, spending history, "
             "employment. The model, optimizing accuracy, INDIRECTLY "
             "reconstructs the forbidden attribute through a proxy — without "
             "any intent.",
             size=12, color=DEEP, line_spacing=1.18)
    teal_callout(s, rx, ly + 2.48, rw, 1.10,
                 "Criterion: a regulated decision requires simultaneously "
                 "reason codes + an appeal path (a human on the contested "
                 "one) + a bias audit of outcomes BEFORE production.",
                 size=12, bold=False)
    text_box(s, rx + 0.04, ly + 3.66, rw, 0.36,
             "The canonical analysis of the mechanism (Obermeyer/Optum) — "
             "Lecture 7.",
             size=10.5, italic=True, color=SLATE)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "«We do not use protected attributes» is a NECESSARY but "
                 "wholly insufficient condition. The only proof is a direct "
                 "audit of outcomes by group.",
                 size=14)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    """case_study — criterion synthesis (automation without gate)."""
    s = blank(p)
    slide_title(s, "The higher the autonomy, the stricter the gate around it "
                   "must be.", size=23, y=0.34, h=0.58, w=12.25)
    bx, by, bw, bh = 0.55, 1.20, 12.25, 1.60
    ocean_box(s, bx, by, bw, bh)
    text_box(s, bx + 0.26, by + 0.14, bw - 0.52, 0.30,
             "One class of error — three different types of AI:", size=14,
             bold=True, color=MID)
    three = [
        ("Zillow", "forecast → automated irreversible buying spree"),
        ("fraud auto-block", "anomaly → automated irreversible block"),
        ("Knight Capital", "deterministic algorithm → automated "
         "irreversible orders"),
    ]
    n = 3
    gap = 0.20
    cwt = (bw - 0.52 - gap * (n - 1)) / n
    tx = bx + 0.26
    for ttl, body in three:
        filled_rect(s, tx, by + 0.52, cwt, 0.94, SURFACE, stroke=SOFT_GREY,
                    stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, tx + 0.12, by + 0.60, cwt - 0.24, 0.30, ttl,
                 size=13, bold=True, color=TEAL)
        text_box(s, tx + 0.12, by + 0.90, cwt - 0.24, 0.50, body,
                 size=11, color=DEEP, line_spacing=1.05)
        tx += cwt + gap
    gold_callout(s, 0.55, 3.00, 12.25, 1.42,
                 "Automation executing irreversible financial actions in an "
                 "open loop without a kill switch (a manual «stop "
                 "everything»), without limits on volume/position, without a "
                 "circuit-breaker (auto-stop on anomaly), and without "
                 "verified deployment, turns an ordinary model error into the "
                 "SPEED OF RUIN.",
                 size=14)
    teal_callout(s, 0.55, 4.46, 12.25, 1.30,
                 "Scoring in Russia is ~100% automated → the criterion "
                 "applies acutely. «100% AI» is acceptable ONLY within a "
                 "harness: reason codes + a human channel + a bias audit + "
                 "drift monitoring + supervision. Remove the harness — and "
                 "it is Apple Card at the scale of a bank.", size=13,
                 bold=False)
    text_box(s, 0.55, 5.92, 12.25, 0.84,
             "High autonomy is neither the goal nor evil in itself. The goal "
             "is autonomy surrounded by paid-for gates, proportional to the "
             "cost of error and the irreversibility of the action.",
             size=13, italic=True, color=MID, line_spacing=1.15)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    build_section_divider(
        p, 4, "LLMs in finance:\nassistants and fact-checking",
        "Until now — deliberately NOT an LLM. Now the LLM — and right away "
        "about its limits in a regulated industry.", "s20")


def build_s21(p):
    """assertion_visual — where LLM fits + fact-integrity base substitution."""
    s = blank(p)
    slide_title(s, "An LLM is good exactly where the task is "
                   "textual-conversational — and bad where it is not.",
                size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.20, 6.35, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.30,
             "Where an LLM is the right type in a bank:", size=13.5,
             bold=True, color=MID)
    fits = [
        "customer support — answers to typical questions",
        "voice assistant — recognize the request, guide through the flow",
        "help for an employee — summarize an inquiry, suggest the "
        "regulation",
        "explaining a product in plain language",
    ]
    fy = ly + 0.52
    for t in fits:
        circle(s, lx + 0.26, fy + 0.05, 0.12, TEAL)
        text_box(s, lx + 0.52, fy, lw - 0.78, 0.56, t,
                 size=12, color=DEEP, line_spacing=1.10)
        fy += 0.58
    text_box(s, lx + 0.26, fy + 0.04, lw - 0.52, 0.34,
             "This is exactly what an LLM is designed for (unlike Sections "
             "1–3).", size=11, italic=True, color=SLATE,
             line_spacing=1.05)
    filled_rect(s, lx + 0.26, fy + 0.46, lw - 0.52, 0.78, SURFACE,
                stroke=SOFT_GREY, stroke_pt=0.75, radius=True,
                radius_adj=0.10)
    text_box(s, lx + 0.42, fy + 0.46, lw - 0.84, 0.78,
             "Scale (verified): T-Bank's chatbot > 40% of inquiries; ~70% of "
             "banks planned voice by 2025.",
             size=11.5, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.08)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 4.05, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 0.18, rw - 0.52, 0.30,
             "Why NOT «> 90% of banks' inquiries»", size=14, bold=True,
             color=TEAL)
    text_box(s, rx + 0.26, ly + 0.58, rw - 0.52, 1.74,
             "The ~90% figure is real — but it is the share of calls handled "
             "by a voice assistant IN ONE BANK'S CALL CENTER, not the share "
             "of banks' inquiries overall. This is BASE SUBSTITUTION "
             "(fact-checking class 5).",
             size=13.5, color=DEEP, line_spacing=1.22)
    text_box(s, rx + 0.26, ly + 2.46, rw - 0.52, 1.50,
             "The lecture teaches fact-checking — building an assertion on a "
             "figure with a substituted base would violate exactly the "
             "principle it teaches. A precise modest figure is worth more "
             "than a loud one.",
             size=12.5, italic=True, color=MID, line_spacing=1.20)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "Symmetry with Sections 1–3: «different task — different "
                 "type of AI» means not «LLM is bad», but «an LLM is good "
                 "exactly where the task is textual-conversational, and bad "
                 "where it is not».",
                 size=14)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """assertion_visual — fact-checking criterion + 5 classes + grounding."""
    s = blank(p)
    slide_title(s, "Here an LLM is not a source of truth, but an interface to "
                   "a source of truth.", size=23, y=0.34, h=0.58, w=12.25)
    # TOP-LEFT — criterion (3 args, compact) ; TOP-RIGHT — 5 error classes
    lx, ly, lw, lh = 0.55, 1.12, 6.10, 2.42
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.24, ly + 0.12, lw - 0.48, 0.28,
             "When an LLM is NOT a source of truth — 3 arguments", size=12.5,
             bold=True, color=MID)
    args = [
        ("mechanism", "generates the plausible, does not retrieve a fact "
                      "(«plausible» ≠ «correct»)"),
        ("cost", "a wrong rate/term = a violation, the organization is "
                 "liable"),
        ("what will break", "free generation about facts = a generator of "
                          "disinformation with legal liability"),
    ]
    ay = ly + 0.46
    for ttl, body in args:
        text_box(s, lx + 0.24, ay, 1.55, 0.62, ttl,
                 size=11.5, bold=True, color=TEAL,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, lx + 1.82, ay, lw - 2.06, 0.62, body,
                 size=10.5, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.08)
        ay += 0.64
    rx, rw = 6.80, 6.00
    ocean_box(s, rx, ly, rw, lh)
    text_box(s, rx + 0.24, ly + 0.12, rw - 0.48, 0.28,
             "Five classes of error in an AI claim", size=12.5, bold=True,
             color=MID)
    classes = [
        "1. hallucinated fact → check against the primary source",
        "2. outdated data → check the date/currency",
        "3. proxy bias in the output → audit outcomes, not inputs",
        "4. deception by metric → FP/FN separately in money",
        "5. base substitution → ask «share of what out of what»",
    ]
    cy = ly + 0.44
    for c in classes:
        text_box(s, rx + 0.26, cy, rw - 0.50, 0.38, c,
                 size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        cy += 0.38
    # BOTTOM BAND — grounding analogy as visual anchor (chapter §4.3)
    gx, gy, gw, gh = 0.55, 3.70, 12.25, 1.62
    ocean_box(s, gx, gy, gw, gh)
    text_box(s, gx + 0.26, gy + 0.12, gw - 0.52, 0.26,
             "Grounding anchor analogy: a student guesses in a confident "
             "tone vs first opens the reference book",
             size=12, bold=True, color=MID)
    add_image(s, DIAGRAMS / "d22-grounding-student.png",
              gx + 0.30, gy + 0.42, gw - 0.60, gh - 0.54)
    gold_callout(s, 0.55, 5.52, 12.25, 1.24,
                 "Alternative: a fixed fact → deterministic retrieval "
                 "(grounded RAG), not generation. In the lecture — recognize "
                 "the class (Understand); verify 5 claims against primary "
                 "sources — Seminar 5 (Apply).", size=14)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    """case_study — Air Canada callback + Klarna arc."""
    s = blank(p)
    slide_title(s, "The right type of AI without grounding and a human exit "
                   "= Air Canada / Klarna.", size=21, y=0.34, h=0.58,
                w=12.25)
    bx, by, bw, bh = 0.55, 1.18, 12.25, 3.30
    ocean_box(s, bx, by, bw, bh)
    col_w = 6.0
    panels = [
        ("Case E — Air Canada (callback)",
         ["the chatbot stated a nonexistent refund policy",
          "the passenger acted on the answer",
          "Feb 2024: the court — the company IS liable for the info of its "
          "chatbot, ordered to pay compensation",
          "class: hallucinated financial fact = legal liability"]),
        ("Case F — Klarna (the 2023 → 2025 arc)",
         ["LLM assistant: ~2/3 of inquiries, ~11 min → < 2 min",
          "claimed savings ~$40M/year (Klarna, 2024)",
          "2024: framed as «AI replaces support»",
          "mid-2025: CSAT↓ → returned to hiring people"]),
    ]
    for j, (ttl, items) in enumerate(panels):
        px = bx + 0.22 + j * (col_w + 0.10)
        filled_rect(s, px, by + 0.18, col_w - 0.20, 0.46,
                    MID if j == 0 else TEAL)
        text_box(s, px, by + 0.18, col_w - 0.20, 0.46, ttl,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        iy = by + 0.78
        for it in items:
            circle(s, px + 0.10, iy + 0.06, 0.11,
                   MID if j == 0 else TEAL)
            text_box(s, px + 0.32, iy, col_w - 0.56, 0.56, it,
                     size=11.5, color=DEEP, line_spacing=1.08)
            iy += 0.60
    teal_callout(s, 0.55, 4.62, 12.25, 0.80,
                 "Criterion: a fixed policy/tariff/right → deterministic "
                 "retrieval/grounded, not free generation; and in any "
                 "scenario a path to a human is guaranteed.", size=13,
                 bold=False)
    gold_callout(s, 0.55, 5.56, 12.25, 1.20,
                 "What is sustainable is not replacement, but AUGMENTATION: "
                 "AI on the mass routine + guaranteed human escalation on the "
                 "emotional / disputed / non-standard tail.",
                 size=15)
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    """assertion_visual — PIVOT checkpoint Р4→Р5 (NOT retro-summary;
    distinct gold angle from s30 payoff; two-level insight = chapter §4.5).
    """
    s = blank(p)
    slide_title(s, "Checkpoint: 4 of 5 types covered — a turn toward «human + "
                   "harness».", size=22, y=0.34, h=0.58, w=12.25)
    # progress checkpoint — 4 done, 1 to go (explicit, not retro-summary)
    bx, by, bw, bh = 0.55, 1.18, 12.25, 1.46
    ocean_box(s, bx, by, bw, bh)
    text_box(s, bx + 0.26, by + 0.12, bw - 0.52, 0.28,
             "Progress across the 5 types — for all 4 covered, a failure of "
             "one form:",
             size=13.5, bold=True, color=MID)
    prog = [("✓ Forecast", "Zillow", TEAL),
            ("✓ Anomalies", "fraud-FP", TEAL),
            ("✓ Scoring", "Apple Card", TEAL),
            ("✓ LLM", "Air Canada", TEAL),
            ("→ Recsys", "next", GOLD)]
    n = 5
    gap = 0.16
    cwf = (bw - 0.52 - gap * (n - 1)) / n
    fx = bx + 0.26
    for ttl, body, col in prog:
        nxt = (col == GOLD)
        filled_rect(s, fx, by + 0.48, cwf, 0.82,
                    GOLD_TINT if nxt else SURFACE,
                    stroke=(GOLD if nxt else SOFT_GREY),
                    stroke_pt=(1.5 if nxt else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, fx + 0.06, by + 0.54, cwf - 0.12, 0.30, ttl,
                 size=12.5, bold=True, color=col, align=PP_ALIGN.CENTER)
        text_box(s, fx + 0.06, by + 0.86, cwf - 0.12, 0.36, body,
                 size=10, color=DEEP, align=PP_ALIGN.CENTER,
                 line_spacing=1.0)
        fx += cwf + gap
    # two-level insight (chapter §4.5 — kept; it is the pivot's substance)
    teal_callout(s, 0.55, 2.80, 12.25, 1.16,
                 "Level 1 «which type of AI» — resolved by the structure of "
                 "the task, often unambiguously. Level 2 «what surrounds the "
                 "AI at the cost of error» — is designed separately. The "
                 "lecture's turn: from here on, attention moves from the type "
                 "to the harness.", size=13.5, bold=True)
    # the pivot itself — why the harmless-looking type is next
    filled_rect(s, 0.55, 4.12, 12.25, 1.16, SURFACE, stroke=LIGHT,
                stroke_pt=1.5, radius=True, radius_adj=0.07)
    icon(s, "shopping-cart", 0.85, 4.46, 0.48, "mid")
    text_box(s, 1.55, 4.22, 11.00, 0.98,
             "The last type — recsys / pricing — seems the most HARMLESS: "
             "the cost of error is near zero. That is exactly why it is more "
             "dangerous than the rest: the failure is SILENT — the system "
             "reports success by its own metric.", size=13, color=DEEP,
             line_spacing=1.16,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.46, 12.25, 1.30,
                 "The question for the remaining section is NOT «which type», "
                 "but «why can the type most harmless by cost of error turn "
                 "out to be more dangerous than the loud failures?»", size=15)
    speaker_notes(s, load_notes("s24"))


def build_s25(p):
    build_section_divider(
        p, 5, "Recommendations and\ndynamic pricing",
        "The last type — the most «harmless» by cost of error, but with the "
        "subtlest pathologies.", "s25")


def build_s26(p):
    """comparison — collaborative vs content-based + user×item matrix."""
    s = blank(p)
    slide_title(s, "Two basic approaches to recommendations — and each with a "
                   "weakness by name.", size=23, y=0.34, h=0.58, w=12.25)
    # TOP — comparison table (full width, compact, parallel structure)
    bx, by, bw, bh = 0.55, 1.12, 12.25, 1.90
    ocean_box(s, bx, by, bw, bh)
    col_w = (bw - 0.40) / 2
    heads = ["Collaborative — «ask those similar to you»",
             "Content-based — «more like it, by description»"]
    for j, hd in enumerate(heads):
        px = bx + 0.20 + j * col_w
        filled_rect(s, px, by + 0.16, col_w - 0.10, 0.42,
                    MID if j == 0 else TEAL)
        text_box(s, px, by + 0.16, col_w - 0.10, 0.42, hd,
                 size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("user × item matrix", "product attributes (genre, brand)"),
        ("⊕ catches unexpected connections", "⊕ no cold-start for a new "
         "product"),
        ("⊖ cold-start + popularity bias", "⊖ over-specialization (niche)"),
    ]
    yy = by + 0.62
    for a, b in rows:
        for j, cc in enumerate((a, b)):
            px = bx + 0.20 + j * col_w
            filled_rect(s, px, yy, col_w - 0.10, 0.40, SURFACE,
                        stroke=SOFT_GREY, stroke_pt=0.75, radius=True,
                        radius_adj=0.14)
            text_box(s, px + 0.18, yy, col_w - 0.36, 0.40, cc,
                     size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.LEFT, line_spacing=1.0)
        yy += 0.42
    # BOTTOM BAND — three sellers analogy as the main teaching visual
    gx, gy, gw, gh = 0.55, 3.16, 12.25, 1.78
    ocean_box(s, gx, gy, gw, gh)
    text_box(s, gx + 0.26, gy + 0.10, gw - 0.52, 0.26,
             "Anchor analogy: three sellers (collaborative · content · "
             "hybrid)", size=12, bold=True, color=MID)
    add_image(s, DIAGRAMS / "d26b-three-sellers.png",
              gx + 0.30, gy + 0.40, gw - 0.60, gh - 0.52)
    gold_callout(s, 0.55, 5.12, 12.25, 0.92,
                 "Remember: WHAT these approaches are and WHICH weakness by "
                 "name each has — collaborative (behavior pattern) vs content "
                 "(the essence of the product).", size=14)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """assertion_visual — hybrid + chart + filter bubble + dynamic pricing."""
    s = blank(p)
    slide_title(s, "A hybrid softens the weaknesses — but without explicit "
                   "diversity it does not cancel the bubble.", size=22,
                y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 6.35, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.14, lw - 0.52, 0.28,
             "Hybrid recommender", size=14, bold=True, color=MID)
    text_box(s, lx + 0.26, ly + 0.46, lw - 0.52, 0.78,
             "collaborative + content-based + context. The content part "
             "closes cold-start, the collaborative one breaks "
             "over-specialization. Most production systems.",
             size=11.5, color=DEEP, line_spacing=1.14)
    # 2 stat-plates instead of ghost bars — readable + honest (estimate)
    pgap = 0.20
    pcw = (lw - 0.52 - pgap) / 2
    for i, (co, val, lab) in enumerate([
            ("Amazon", "~35%", "of revenue — from recommendations"),
            ("Netflix", "~75%", "of viewing — from recommendations")]):
        ppx = lx + 0.26 + i * (pcw + pgap)
        filled_rect(s, ppx, ly + 1.34, pcw, 1.46, SURFACE, stroke=LIGHT,
                    stroke_pt=1.5, radius=True, radius_adj=0.10)
        text_box(s, ppx + 0.10, ly + 1.44, pcw - 0.20, 0.28, co,
                 size=12.5, bold=True, color=MID, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, ppx + 0.06, ly + 1.74, pcw - 0.12, 0.62, val,
                 size=30, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, ppx + 0.10, ly + 2.36, pcw - 0.20, 0.38, lab,
                 size=10.5, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, lx + 0.26, ly + 2.92, lw - 0.52, 0.98,
             "A historical estimate (dating back to McKinsey ~2013), NOT a "
             "fresh verified figure — hence a number, not a chart-«fact».",
             size=10.5, italic=True, color=SLATE, line_spacing=1.12)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 1.92)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.28,
             "Filter bubble", size=13, bold=True,
             color=MID)
    text_box(s, rx + 0.24, ly + 0.46, rw - 0.48, 1.36,
             "showed something similar → the user interacts with the similar "
             "→ the model shows something even more similar. The narrowing of "
             "diversity is NOT malicious intent, but a consequence of "
             "optimizing for short-term relevance.",
             size=11.5, color=DEEP, line_spacing=1.16)
    ocean_box(s, rx, ly + 2.08, rw, 1.97)
    text_box(s, rx + 0.24, ly + 2.22, rw - 0.48, 0.28,
             "Dynamic pricing", size=13, bold=True, color=MID)
    text_box(s, rx + 0.24, ly + 2.54, rw - 0.48, 1.42,
             "auto-adjustment of price to demand/time/context. The "
             "perception of fairness and the law are CONSTRAINTS of the "
             "task, NOT variables to optimize.",
             size=11.5, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "The filter bubble arises without any intent — a direct "
                 "consequence of «guess what you'll buy now». A hybrid does "
                 "NOT automatically cancel this narrowing.", size=14)
    footer(s, "Amazon ~35% / Netflix ~75% — a historical estimate (McKinsey "
              "~2013), NOT a fresh figure; Ozon/WB — the companies do not "
              "disclose the exact share.")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """case_study — proxy≠goal + Wendy's + criterion."""
    s = blank(p)
    slide_title(s, "The system reports success by its own metric — while "
                   "destroying what did not make it into the metric.",
                size=21,
                y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 7.05, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.14, lw - 0.52, 0.28,
             "The root — proxy ≠ goal", size=14, bold=True, color=MID)
    text_box(s, lx + 0.26, ly + 0.46, lw - 0.52, 0.94,
             "the system optimizes what is measurable in the moment (the "
             "PROXY: clicks, time, margin) ≠ the real goal (trust, "
             "well-being). The divergence → filter bubble, dark patterns, "
             "homogenization.",
             size=11.5, color=DEEP, line_spacing=1.16)
    filled_rect(s, lx + 0.26, ly + 1.50, lw - 0.52, 0.34, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.2, radius=True, radius_adj=0.2)
    text_box(s, lx + 0.40, ly + 1.50, lw - 0.80, 0.34,
             "Case G — Wendy's, February 2024", size=12.5, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 1.94, lw - 0.52, 1.96,
             "~$20M into digital menu boards with dynamic pricing → the "
             "media read it as surge-pricing → #BoycottWendys, a competitor "
             "outplayed it in advertising → rollback within days. The money "
             "was invested, the technology was available — what failed was "
             "ignoring that perceived fairness = a CONSTRAINT, not a "
             "variable.",
             size=11.5, color=DEEP, line_spacing=1.18)
    rx, rw = 7.80, 5.00
    # counter-weight: «прокси ≠ цель» mini-diagram (chapter §5.5, derived)
    ocean_box(s, rx, ly, rw, 1.94)
    text_box(s, rx + 0.22, ly + 0.12, rw - 0.44, 0.26,
             "Why the failure is SILENT:", size=12.5, bold=True, color=MID)
    # rising proxy line vs flat/declining true goal
    gx0, gy0, gx1 = rx + 0.30, ly + 1.62, rx + rw - 0.30
    connector(s, gx0, gy0, gx0, ly + 0.52, color=SLATE, width=1.2)
    connector(s, gx0, gy0, gx1, gy0, color=SLATE, width=1.2)
    connector(s, gx0, ly + 1.42, gx1, ly + 0.62, color=GOLD, width=3.0)
    connector(s, gx0, ly + 1.12, gx1, ly + 1.30, color=TEAL, width=3.0,
              dash="dash")
    text_box(s, gx1 - 1.70, ly + 0.50, 1.70, 0.26, "proxy metric ↑",
             size=10, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    text_box(s, gx1 - 1.70, ly + 1.30, 1.70, 0.26, "true goal →",
             size=10, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    text_box(s, rx + 0.22, ly + 1.66, rw - 0.44, 0.24,
             "dashboard green, trust falling", size=10,
             italic=True, color=SLATE)
    teal_callout(s, rx, ly + 2.10, rw, 1.95,
                 "Criterion:\n• proxy ≠ goal → serendipity + explainability "
                 "+ an audit for discrimination\n• the cost is a decision of "
                 "a HUMAN IN A LEGAL FRAME, not the output of the "
                 "optimizer", size=12,
                 bold=False)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "One class runs through the whole lecture: accuracy "
                 "deception, proxy bias in scoring, Klarna, filter bubble — "
                 "PROXY ≠ GOAL. The more powerful the optimizer, the costlier "
                 "the divergence. The failure is silent.", size=14)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """matrix / schema_matrix — task × AI-type, 6 rows."""
    s = blank(p)
    slide_title(s, "Name the type for the structure of the task — and "
                   "sometimes the right type is not AI.", size=22, y=0.32,
                h=0.56, w=12.25)
    bx, by, bw, bh = 0.40, 1.00, 12.55, 4.80
    ocean_box(s, bx, by, bw, bh)
    tx, ty = bx + 0.16, by + 0.14
    col_w = [3.05, 4.95, 4.23]
    headers = ["Task", "Type of AI · why that one", "Typical failure"]
    hh = 0.46
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID)
        text_box(s, cx + 0.10, ty, col_w[j] - 0.18, hh, hd,
                 size=12.5, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE,
                 align=(PP_ALIGN.LEFT if j else PP_ALIGN.LEFT))
        cx += col_w[j]
    rows = [
        ("trending-up", "Demand forecasting / cash-flow / churn",
         "Time-series forecasting (ARIMA/boosting) — numbers with a trend, "
         "not text",
         "distribution shift × irreversible automation (Zillow)", False),
        ("radar", "Real-time fraud / AML",
         "Anomaly detection + rules-engine — norm/deviation",
         "FP at scale; accuracy lies under imbalance", False),
        ("scale", "Credit scoring",
         "Tabular ML (logreg/GBM+SHAP) — explainability = law",
         "proxy bias + opacity (Apple Card)", False),
        ("message-circle", "Support / voice / explanation",
         "LLM (grounded) — a textual-conversational task",
         "hallucinated financial fact = legal liability (Air Canada)",
         False),
        ("shopping-cart", "Recommendations / pricing",
         "Recsys; pricing — an optimizer within a frame",
         "proxy≠goal: filter bubble (Wendy's)", False),
        ("circle-slash", "A deterministic regulatory task (a hard AML "
         "threshold)",
         "Ordinary code / rules-engine — NOT AI; law ≠ probability",
         "AI would add non-determinism + an error surface", True),
    ]
    rh = 0.62
    yy = ty + hh
    for ic, task, aitype, fail, hi in rows:
        cx = tx
        cells = [task, aitype, fail]
        for j, cc in enumerate(cells):
            bg = GOLD_TINT if hi else (
                WHITE if (rows.index(
                    (ic, task, aitype, fail, hi)) % 2 == 0) else SURFACE)
            filled_rect(s, cx, yy, col_w[j], rh, bg,
                        stroke=(GOLD if hi else SOFT_GREY),
                        stroke_pt=(1.3 if hi else 0.75))
            if j == 0:
                icon(s, ic, cx + 0.10, yy + (rh - 0.30) / 2, 0.30,
                     "gold" if hi else "mid")
                text_box(s, cx + 0.48, yy, col_w[j] - 0.56, rh, cc,
                         size=10.5, bold=True,
                         color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                         line_spacing=1.0)
            else:
                text_box(s, cx + 0.12, yy, col_w[j] - 0.22, rh, cc,
                         size=10.5, bold=(hi and j == 1),
                         color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                         line_spacing=1.02)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 5.92, 12.25, 0.86,
                 "The bottom row — deliberately. A hard, verifiable "
                 "regulatory rule → deterministic code, NOT AI: AI would add "
                 "non-determinism where accuracy and auditability are "
                 "needed.", size=14)
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """assertion_visual — FINAL PAYOFF: one principle, distinct from s29
    matrix (operational apparatus) and s24 pivot. Dominant principle
    statement on top; 5 cases as compact evidence strip below."""
    s = blank(p)
    slide_title(s, "The substantive answer to the central question of the "
                   "lecture.",
                size=23, y=0.34, h=0.56, w=12.25)
    # DOMINANT principle band — the single payoff statement
    bx, by, bw, bh = 0.55, 1.06, 12.25, 1.84
    ocean_box(s, bx, by, bw, bh, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, bx + 0.50, by + 0.20, bw - 1.00, bh - 0.40,
             "In none of the five failures is the fix «a better model». In "
             "all of them — a better JUDGMENT about what stands around the "
             "AI at the cost of error.",
             size=21, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, line_spacing=1.20)
    # 5 cases as a compact horizontal evidence strip (case → fix)
    text_box(s, 0.55, 3.06, 12.25, 0.28,
             "Five cases from the lecture — one and the same remedy (NOT a "
             "model):",
             size=12.5, bold=True, color=MID)
    strip = [
        ("Zillow", "human-gate"),
        ("Apple Card", "explainability + audit"),
        ("Air Canada", "grounding"),
        ("Klarna", "augmentation"),
        ("Wendy's", "legal frame"),
    ]
    n = 5
    gap = 0.18
    cw = (12.25 - gap * (n - 1)) / n
    sx, sy, sh = 0.55, 3.42, 1.46
    for case, fix in strip:
        ocean_box(s, sx, sy, cw, sh)
        text_box(s, sx + 0.08, sy + 0.16, cw - 0.16, 0.34, case,
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, sx + cw / 2 - 0.10, sy + 0.54, 0.20, 0.26, "↓",
                 size=15, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
        filled_rect(s, sx + 0.14, sy + 0.82, cw - 0.28, 0.50, TEAL_TINT,
                    stroke=TEAL, stroke_pt=1.0, radius=True,
                    radius_adj=0.16)
        text_box(s, sx + 0.18, sy + 0.82, cw - 0.36, 0.50, fix,
                 size=10.5, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        sx += cw + gap
    gold_callout(s, 0.55, 5.06, 12.25, 1.00,
                 "This is the answer: not «AI is good/bad», but «name the "
                 "type, justify the choice, and design the harness for the "
                 "cost of error».", size=14)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """assertion_visual — security 2 panels (А law / Б CV) — AGGREGATED."""
    s = blank(p)
    slide_title(s, "Financial data, PII (personal data), and biometrics must "
                   "not go into a public LLM.", size=21, y=0.32, h=0.56,
                w=12.25)
    # TWO clearly separated panels with header plates + thick divider
    lx, ly, lw, lh = 0.55, 1.06, 6.00, 2.78
    ocean_box(s, lx, ly, lw, lh)
    filled_rect(s, lx, ly, lw, 0.50, MID, radius=True, radius_adj=0.06)
    icon(s, "lock", lx + 0.22, ly + 0.07, 0.36, "white")
    text_box(s, lx + 0.70, ly, lw - 0.90, 0.50,
             "(A) Data and the law: 152-FZ / PII", size=13, bold=True,
             color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    a_pts = [
        "financial data + PII + biometrics = sensitive; Federal Law 152-FZ "
        "— localization of personal data of Russian citizens, biometrics = a "
        "strict regime",
        "public cloud: the data leaves the loop · no control of retention · "
        "auditability drops",
        "the criterion is by data sensitivity and regime, NOT by the "
        "strength of the model",
    ]
    ay = ly + 0.62
    for t in a_pts:
        circle(s, lx + 0.24, ay + 0.05, 0.11, MID)
        text_box(s, lx + 0.48, ay, lw - 0.72, 0.70, t,
                 size=11, color=DEEP, line_spacing=1.10)
        ay += 0.72
    # vertical divider between the two independent sub-topics
    filled_rect(s, 6.62, ly, 0.06, lh, LIGHT)
    rx, rw = 6.80, 6.00
    ocean_box(s, rx, ly, rw, lh)
    filled_rect(s, rx, ly, rw, 0.50, TEAL, radius=True, radius_adj=0.06)
    icon(s, "scan-face", rx + 0.22, ly + 0.07, 0.36, "white")
    text_box(s, rx + 0.70, ly, rw - 0.90, 0.50,
             "(B) The CV layer: KYC, biometrics, hidden labor", size=13,
             bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    b_pts = [
        "KYC = identifying the customer; liveness — a live human in front of "
        "the camera, not a photo/mask/deepfake (a strict regime)",
        "Just Walk Out (Amazon) shut down in 2024: the «autonomous» checkout "
        "relied on > 1000 reviewers in India — hidden labor",
        "Computer-vision bias is deepened in Lecture 7.",
    ]
    by2 = ly + 0.62
    for t in b_pts:
        circle(s, rx + 0.24, by2 + 0.05, 0.11, TEAL)
        text_box(s, rx + 0.48, by2, rw - 0.72, 0.70, t,
                 size=11, color=DEEP, line_spacing=1.10)
        by2 += 0.72
    # GOLD ANCHOR BAND — biometrics irreversible (password vs face) — large
    gx, gy, gw, gh = 0.55, 3.96, 12.25, 1.98
    filled_rect(s, gx, gy, gw, gh, GOLD_TINT, stroke=GOLD, stroke_pt=2.5,
                radius=True, radius_adj=0.05)
    text_box(s, gx + 0.34, gy + 0.16, 3.95, gh - 0.32,
             "Biometrics are irreversibly\ncompromised on a\nleak — the same\n"
             "logic «irreversible →\na strict gate» that\nran through the "
             "whole\nlecture.",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16)
    add_image(s, DIAGRAMS / "d31-password-vs-face.png",
              gx + 4.40, gy + 0.14, gw - 4.70, gh - 0.28)
    gold_callout(s, 0.55, 6.08, 12.25, 0.78,
                 "«Fully autonomous» in marketing is a hypothesis to be "
                 "tested, not a fact (Amazon disputed the scale).",
                 size=14)
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    """qa_minimal — checklist + Семинар 5 bridge + Q&A."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    bx, by, bw, bh = 0.55, 0.45, 12.25, 2.70
    ocean_box(s, bx, by, bw, bh)
    text_box(s, bx + 0.26, by + 0.14, bw - 0.52, 0.30,
             "Before applying AI to a financial / retail task:",
             size=14, bold=True, color=MID)
    checks = [
        "1. Which type of AI does the structure of the task dictate?",
        "2. Why NOT an LLM here? (≥1 structural reason)",
        "3. Can it be solved without AI at all? (rule → code)",
        "4. Is the action at the output reversible? (irreversible → gate)",
        "5. How to verify the fact? (class of error + principle)",
        "6. Is the decision regulated? (reason codes / appeal)",
        "7. Who is liable and where is a human mandatory?",
        "8. Is there PII / financial data / biometrics? → not into an LLM; "
        "152-FZ",
    ]
    n = len(checks)
    col_w = (bw - 0.52) / 2
    for i, c in enumerate(checks):
        col = i % 2
        row = i // 2
        cx = bx + 0.26 + col * col_w
        cy = by + 0.52 + row * 0.52
        text_box(s, cx, cy, col_w - 0.16, 0.48, c,
                 size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    ocean_box(s, bx, 3.30, bw, 1.30, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, bx + 0.30, 3.42, bw - 0.6, 0.34,
             "Assignment — Seminar 5 «Fact-checking AI on financial data»",
             size=14, bold=True, color=DEEP)
    text_box(s, bx + 0.30, 3.78, bw - 0.6, 0.72,
             "In ~30 min, teams verify 5 real claims against primary sources "
             "themselves (Bank of Russia / VCIOM), marking them "
             "true/partial/false. The lecture teaches recognizing the class "
             "(Understand); the seminar — verifying independently (Apply).",
             size=12, color=DEEP, line_spacing=1.16)
    text_box(s, 0.55, 4.78, 12.25, 1.45, "Questions",
             size=92, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.55, 6.30, 12.25, 0.50,
             "The palette of types — a lens for all the industry lectures "
             "ahead.  Thank you for your attention.", size=15, italic=True,
             color=LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s32"))


# ============================================================
# Main
# ============================================================
def main():
    spec = load_deck()
    if spec is not None:
        n = len(spec["slides"])
        print(f"deck spec OK — {n} slides (deck.yaml + deck-part2.yaml), "
              f"totals {spec['totals'].get('slides')}")
    p = setup_pres()
    builders = [
        build_s01, build_s02, build_s03, build_s04, build_s04a,
        build_s05, build_s06, build_s07, build_s08, build_s09,
        build_s10, build_s11, build_s12, build_s13, build_s14,
        build_s15, build_s16, build_s17, build_s18, build_s19,
        build_s20, build_s21, build_s22, build_s23, build_s24,
        build_s25, build_s26, build_s27, build_s28, build_s29,
        build_s30, build_s31, build_s32]
    assert len(builders) == 33, f"expected 33 builders, got {len(builders)}"
    for b in builders:
        b(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"saved {OUT} — {len(p.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
