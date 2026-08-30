#!/usr/bin/env python3
"""notes_pages_pdf.py — reusable «notes-pages PDF» builder for course lectures.

Produces a PORTRAIT (A4) handout PDF where every page shows one lecture slide
(rendered image) on top and that slide's speaker notes (readable, word-wrapped,
Cyrillic) below. Long notes are NEVER clipped — they overflow onto continuation
pages that carry the SAME compact header (no image), so nothing is truncated.

Per-page layout (owner spec):
  ┌───────────────────────────────────────────────────────────────┐
  │ HEADER (muted, ~9-10pt): «‹Lecture title› · ‹Slide title› ·    │
  │                           слайд N»                             │
  │ ┌─────────────── slide image (16:9, full width) ────────────┐ │
  │ │                                                           │ │
  │ └───────────────────────────────────────────────────────────┘ │
  │ ───────────────────────────────────────────────────────────── │
  │ speaker notes … narrative + inline [N] + «Источники:» URLs …  │
  │                                                               │
  │                                    FOOTER (right): «P / TOTAL» │
  └───────────────────────────────────────────────────────────────┘

DATA SOURCES (key — everything read from the RENDERED artefacts, not .md):
  • Slide images  — pages of `rendered/lec-NN.pdf`, IN ORDER (page i == slide i,
    already carrying page numbers + refs). Rasterised at --dpi.
  • Notes         — `rendered/lec-NN.pptx` via python-pptx: for each slide in
    presentation order, `slide.notes_slide.notes_text_frame.text`. This gives
    the FULL baked note — narrative + inline [N] markers + the trailing
    «Источники:» URL block — which lives ONLY in the pptx (appended at build
    time from the ref registry, never written back to slides/*.md). Reading the
    pptx (a) restores the reference URLs into the notes-PDF, and (b) makes
    slide↔note matching purely POSITIONAL (PDF page i ↔ pptx slide i, both in
    presentation order) — no fragile filename / natural-sort / BUILD_ORDER
    logic. Decks with letter-suffix slides or non-monotone build order can no
    longer mis-pair notes with images.
  • Lecture title — `deck.yaml` → `deck.title` (full, unabbreviated).
  • Slide title   — the slide's title placeholder if present; otherwise the
    topmost / largest text block of the slide (its headline first line). Full &
    readable — never an «S01» abbreviation.

OUTPUT: rendered/lec-NN-notes.pdf

USAGE:
  python3 tools/presentation-build/notes_pages_pdf.py library/lectures/lec-04
  python3 tools/presentation-build/notes_pages_pdf.py library/lectures/lec-01 --dpi 150

Parameterised by lecture folder so it applies unchanged to lec-01 … lec-NN.
Requires: pymupdf (fitz) + python-pptx. Cyrillic via a bundled TTF (DejaVu /
Liberation / Noto — auto-discovered; override with --font).
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

import pymupdf  # PyMuPDF ≥ 1.24
from pptx import Presentation  # python-pptx ≥ 0.6
from pptx.util import Emu


# --- A4 portrait geometry (points; 1pt = 1/72") ---
PAGE_W = 595.0
PAGE_H = 842.0
MARGIN = 42.0                      # outer margin
GUTTER = 16.0                      # gap between slide image and notes

# --- typography (points) ---
NOTE_SIZE = 11.5
NOTE_LEADING = 15.5                # line pitch
HEAD_SIZE = 9.5                    # muted header (~9-10pt per owner spec)
FOOT_SIZE = 9.0

# --- colours (0-1 RGB) ---
INK = (0.13, 0.16, 0.36)          # DEEP #21295C
MUTED = (0.42, 0.47, 0.55)        # header — muted slate-blue
SLATE = (0.36, 0.40, 0.47)
RULE = (0.83, 0.87, 0.92)
FRAME = (0.11, 0.45, 0.58)        # LIGHT #1C7293

# --- font auto-discovery (Cyrillic-capable) ---
_FONT_CANDIDATES = [
    ("DejaVuSans.ttf", "DejaVuSans-Bold.ttf"),
    ("LiberationSans-Regular.ttf", "LiberationSans-Bold.ttf"),
    ("NotoSans-Regular.ttf", "NotoSans-Bold.ttf"),
    ("Arial.ttf", "Arial Bold.ttf"),
]
_FONT_SEARCH_DIRS = [
    "/home/harness/.local/lo-sysroot/usr/share/fonts",
    "/home/harness/.local/libreoffice-portable/share/fonts",
    "/tmp/claude-999/local/usr/share/fonts",
    "/usr/share/fonts",
    "/usr/local/share/fonts",
    str(Path.home() / ".fonts"),
    str(Path.home() / ".local/share/fonts"),
]


def _find_font(regular_override: str | None = None):
    """Return (regular_path, bold_path). bold_path may be None (falls back to
    regular). Prefers a full Cyrillic-capable family."""
    if regular_override:
        rp = Path(regular_override)
        if rp.exists():
            bold = None
            for cand in (rp.with_name(rp.stem + "-Bold" + rp.suffix),
                         rp.with_name(rp.stem + " Bold" + rp.suffix)):
                if cand.exists():
                    bold = str(cand)
                    break
            return str(rp), bold
        print(f"[warn] --font {regular_override} not found; auto-discovering",
              file=sys.stderr)

    found: dict[str, str] = {}
    for base in _FONT_SEARCH_DIRS:
        bd = Path(base)
        if not bd.exists():
            continue
        for f in bd.rglob("*.ttf"):
            found.setdefault(f.name, str(f))
    for reg, bold in _FONT_CANDIDATES:
        if reg in found:
            return found[reg], found.get(bold)
    raise SystemExit(
        "No Cyrillic-capable TTF found. Pass --font /path/to/Font.ttf "
        f"(searched {_FONT_SEARCH_DIRS}).")


# ------------------------------------------------------------------
# Lecture title — from deck.yaml (deck.title).
# ------------------------------------------------------------------
def lecture_title(lecture_dir: Path) -> str:
    """Full lecture title from deck.yaml `deck.title`. Minimal YAML scan (avoid
    a PyYAML dependency): find the `deck:` block, then its `title:` key."""
    for name in ("deck.yaml",):
        p = lecture_dir / name
        if not p.exists():
            continue
        in_deck = False
        for raw in p.read_text(encoding="utf-8").splitlines():
            line = raw.rstrip("\n")
            if re.match(r'^\s*#', line) or not line.strip():
                continue
            if re.match(r'^deck:\s*$', line):
                in_deck = True
                continue
            if in_deck:
                # leaving the deck: block when a new top-level key appears
                if re.match(r'^\S', line) and not re.match(r'^deck:', line):
                    in_deck = False
                    continue
                m = re.match(r'^\s+title:\s*(.+?)\s*$', line)
                if m:
                    return m.group(1).strip().strip('"').strip("'")
    return lecture_dir.name


# ------------------------------------------------------------------
# Slide title — placeholder if present, else topmost/largest text block.
# ------------------------------------------------------------------
_DECOR_RE = re.compile(r'^[\W\d_]{0,4}$')


def _is_decorative(first_line: str) -> bool:
    """Skip pure-number / tiny decorative tokens (e.g. the big «04» on a cover
    slide, page markers) when choosing a slide title."""
    s = first_line.strip()
    if not s:
        return True
    if re.fullmatch(r'\d+', s):
        return True
    if re.fullmatch(r'\d+\s*/\s*\d+', s):
        return True
    # 4 chars or fewer AND no run of 3 letters → decorative glyphs
    if len(s) <= 4 and not re.search(r'[а-яёa-z]{3}', s.lower()):
        return True
    return False


def slide_title(slide) -> str:
    """Return a full, readable slide title.

    Preference: (1) title placeholder text; else (2) among text blocks in the
    top band of the slide, the largest-font block's first line (decorative
    pure-number / tiny tokens skipped)."""
    # (1) title placeholder
    try:
        tph = slide.shapes.title
    except Exception:
        tph = None
    if tph is not None and tph.has_text_frame:
        t = tph.text_frame.text.strip().split("\n")[0].strip()
        if t and not _is_decorative(t):
            return t

    # (2) topmost / largest text block
    cands = []  # (top_in, max_font_pt, first_line)
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        first_line = txt.split("\n")[0].strip()
        if not first_line:
            continue
        top_in = Emu(sh.top).inches if sh.top is not None else 99.0
        sizes = [r.font.size.pt
                 for para in sh.text_frame.paragraphs
                 for r in para.runs
                 if r.font.size is not None]
        max_sz = max(sizes) if sizes else 0.0
        cands.append((top_in, max_sz, first_line))

    if not cands:
        return "(без заголовка)"

    cands.sort(key=lambda c: c[0])
    top_min = cands[0][0]
    band = [c for c in cands if c[0] <= top_min + 1.5]
    real = [c for c in band if not _is_decorative(c[2])]
    pool = real if real else band
    # largest font first, then topmost
    pool.sort(key=lambda c: (-c[1], c[0]))
    return pool[0][2]


# ------------------------------------------------------------------
# Notes extraction — read the rendered deck (pptx) in presentation order.
# ------------------------------------------------------------------
def notes_from_pptx(pptx_path: Path) -> list[str]:
    """Return speaker-note strings, one per slide, in presentation order
    (== page order of rendered/lec-NN.pdf). FULL baked note: narrative + inline
    [N] + «Источники:» URL block. A slide without notes → empty string."""
    prs = Presentation(str(pptx_path))
    out: list[str] = []
    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                text = tf.text or ""
        out.append(text.strip())
    return out


def slide_titles_from_pptx(pptx_path: Path) -> list[str]:
    """Return slide titles, one per slide, in presentation order."""
    prs = Presentation(str(pptx_path))
    return [slide_title(slide) for slide in prs.slides]


def notes_to_paragraphs(notes: str) -> list[str]:
    """Blank-line-separated blocks → paragraphs; single newlines inside a block
    collapse to spaces. A «Источники:» block keeps its per-URL line breaks."""
    if not notes:
        return ["(нет заметок для этого слайда)"]
    blocks = [b.strip() for b in re.split(r'\n\s*\n', notes) if b.strip()]
    out = []
    for b in blocks:
        if b.lstrip().startswith("Источники:"):
            out.extend(ln.rstrip() for ln in b.split("\n") if ln.strip())
        else:
            out.append(re.sub(r'\s*\n\s*', ' ', b))
    return out


# ------------------------------------------------------------------
# Slide-image source discovery (PDF pages, in order).
# ------------------------------------------------------------------
def open_slide_source(rendered: Path, lec_id: str):
    """Return (kind, handle, count). Prefer rendered/lec-NN.pdf (page i = slide
    i+1); fall back to snapshots/slide-*.png sorted."""
    pdf = rendered / f"{lec_id}.pdf"
    if pdf.exists():
        doc = pymupdf.open(pdf)
        if doc.page_count > 0:
            return "pdf", doc, doc.page_count
        doc.close()
    snaps = sorted((rendered / "snapshots").glob("slide-*.png"))
    if snaps:
        return "png", snaps, len(snaps)
    snaps = sorted((rendered / "snapshots").glob("*.png"))
    if snaps:
        return "png", snaps, len(snaps)
    raise SystemExit(
        f"No slide images: neither {pdf} nor {rendered}/snapshots/slide-*.png")


def slide_pixmap(kind, handle, idx0: int, dpi: int) -> pymupdf.Pixmap:
    """idx0 is 0-based slide index."""
    if kind == "pdf":
        return handle[idx0].get_pixmap(dpi=dpi)
    return pymupdf.Pixmap(str(handle[idx0]))


# ------------------------------------------------------------------
# Word-wrap that respects the chosen font's real metrics.
# ------------------------------------------------------------------
def wrap_paragraph(font: pymupdf.Font, text: str, size: float,
                   max_w: float) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines, cur = [], ""
    for w in words:
        trial = w if not cur else cur + " " + w
        if font.text_length(trial, size) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            if font.text_length(w, size) > max_w:
                # a single word (e.g. long URL) longer than the line: hard-break
                piece = ""
                for ch in w:
                    if font.text_length(piece + ch, size) <= max_w:
                        piece += ch
                    else:
                        lines.append(piece)
                        piece = ch
                cur = piece
            else:
                cur = w
    if cur:
        lines.append(cur)
    return lines


def _elide(font: pymupdf.Font, text: str, size: float, max_w: float) -> str:
    """Single-line elide with a trailing «…» to fit max_w (header only)."""
    if font.text_length(text, size) <= max_w:
        return text
    ell = "…"
    out = ""
    for ch in text:
        if font.text_length(out + ch + ell, size) > max_w:
            break
        out += ch
    return (out.rstrip() + ell) if out else ell


# ------------------------------------------------------------------
# Page rendering.
# ------------------------------------------------------------------
def build(lecture_dir: Path, dpi: int, font_override: str | None):
    lec_id = lecture_dir.name                       # e.g. "lec-04"
    rendered = lecture_dir / "rendered"
    if not rendered.exists():
        raise SystemExit(f"missing {rendered}")

    pptx_path = rendered / f"{lec_id}.pptx"
    if not pptx_path.exists():
        raise SystemExit(f"missing {pptx_path}")

    reg_path, bold_path = _find_font(font_override)
    font = pymupdf.Font(fontfile=reg_path)
    FNAME, FBOLD = "notesreg", "notesbold"

    lec_title = lecture_title(lecture_dir)
    notes_all = notes_from_pptx(pptx_path)
    titles_all = slide_titles_from_pptx(pptx_path)

    kind, handle, n_slides = open_slide_source(rendered, lec_id)
    # positional alignment guard: PDF page i ↔ pptx slide i
    n = min(n_slides, len(notes_all), len(titles_all))
    if not (n_slides == len(notes_all) == len(titles_all)):
        print(f"[warn] count mismatch: pdf_pages={n_slides} "
              f"pptx_notes={len(notes_all)} pptx_titles={len(titles_all)}; "
              f"using {n}", file=sys.stderr)

    out = pymupdf.open()

    content_w = PAGE_W - 2 * MARGIN

    def new_page():
        pg = out.new_page(width=PAGE_W, height=PAGE_H)
        pg.insert_font(fontname=FNAME, fontfile=reg_path)
        pg.insert_font(fontname=FBOLD, fontfile=bold_path or reg_path)
        return pg

    def draw_header(pg, lec_t, slide_t, slide_no):
        """Muted header (~9-10pt): «‹lecture title› · ‹slide title› · слайд N».

        All three parts must survive. If the full line is too wide, elide ONLY
        the slide-title middle — the lecture-title prefix and the «· слайд N»
        suffix are always kept whole (owner spec §3: header must carry lecture
        title + slide title + slide number)."""
        y = MARGIN
        sep = " · "
        prefix = lec_t + sep
        suffix = sep + f"слайд {slide_no}"
        full = prefix + slide_t + suffix
        if font.text_length(full, HEAD_SIZE) <= content_w:
            line = full
        else:
            avail = content_w - font.text_length(prefix + suffix, HEAD_SIZE)
            line = prefix + _elide(font, slide_t, HEAD_SIZE, max(avail, 0)) + suffix
        pg.insert_text((MARGIN, y + HEAD_SIZE), line, fontname=FNAME,
                       fontsize=HEAD_SIZE, color=MUTED)
        return y + HEAD_SIZE + 8

    def draw_slide_image(pg, top, idx0):
        pix = slide_pixmap(kind, handle, idx0, dpi)
        iw, ih = pix.width, pix.height
        ratio = iw / ih if ih else 16 / 9
        box_w = content_w
        box_h = box_w / ratio
        rect = pymupdf.Rect(MARGIN, top, MARGIN + box_w, top + box_h)
        pg.draw_rect(rect, color=FRAME, width=0.8)
        pg.insert_image(rect, pixmap=pix)
        return top + box_h

    def draw_notes_rule(pg, top):
        pg.draw_line((MARGIN, top), (PAGE_W - MARGIN, top),
                     color=RULE, width=0.8)
        return top + 12

    def draw_footer(pg, doc_page_no, total_pages):
        """Footer = document page number only (owner spec §1-2): «N / total».
        NO «lec-NN · SNN»."""
        text = f"{doc_page_no} / {total_pages}"
        pg.insert_text((PAGE_W - MARGIN - font.text_length(text, FOOT_SIZE),
                        PAGE_H - MARGIN + 14), text,
                       fontname=FNAME, fontsize=FOOT_SIZE, color=SLATE)

    notes_bottom = PAGE_H - MARGIN - 18   # leave room for footer
    # image box height for slide-image pages (16:9 assumption for layout math;
    # actual per-image ratio is applied at draw time, near-identical for 16:9)
    img_h = content_w * 9.0 / 16.0
    # y at which notes start on a slide-image page vs a continuation page
    y_notes_first = MARGIN + HEAD_SIZE + 8 + img_h + GUTTER + 12
    y_notes_cont = MARGIN + HEAD_SIZE + 8 + 2 + 12

    def line_advance(line: str) -> float:
        return NOTE_LEADING if line else (NOTE_LEADING * 0.55)

    # ---- PASS 1: wrap notes + lay out pages (no rasterisation) ----
    # Each slide → list of "pages": page 0 carries the image, the rest are
    # continuations. A page is a list of (line, is_source) to draw.
    per_slide_pages: list[list[list[tuple[str, bool]]]] = []
    for idx0 in range(n):
        wrapped: list[tuple[str, bool]] = []
        for para in notes_to_paragraphs(notes_all[idx0]):
            is_src = para.startswith("Источники:") or para.startswith("[")
            for ln in wrap_paragraph(font, para, NOTE_SIZE, content_w):
                wrapped.append((ln, is_src))
            wrapped.append(("", is_src))        # paragraph gap
        if wrapped and wrapped[-1][0] == "":
            wrapped.pop()

        pages: list[list[tuple[str, bool]]] = [[]]
        y = y_notes_first
        for line, is_src in wrapped:
            if y + NOTE_LEADING > notes_bottom:
                pages.append([])
                y = y_notes_cont
            pages[-1].append((line, is_src))
            y += line_advance(line)
        per_slide_pages.append(pages)

    total_pages = sum(len(p) for p in per_slide_pages)

    # ---- PASS 2: render, stamping the footer inline on each page ----
    doc_page_no = 0
    for idx0 in range(n):
        pages = per_slide_pages[idx0]
        s_title = titles_all[idx0]
        slide_no = idx0 + 1
        for p_i, page_lines in enumerate(pages):
            doc_page_no += 1
            pg = new_page()
            y = draw_header(pg, lec_title, s_title, slide_no)
            if p_i == 0:
                # slide-image page
                y = draw_slide_image(pg, y, idx0) + GUTTER
                y = draw_notes_rule(pg, y)
            else:
                # continuation page: SAME header (no «продолжение»), NO image
                y = draw_notes_rule(pg, y + 2)
            for line, is_src in page_lines:
                if line:
                    col = SLATE if is_src else INK
                    sz = NOTE_SIZE - 1 if is_src else NOTE_SIZE
                    pg.insert_text((MARGIN, y), line, fontname=FNAME,
                                   fontsize=sz, color=col)
                y += line_advance(line)
            draw_footer(pg, doc_page_no, total_pages)

    out_path = rendered / f"{lec_id}-notes.pdf"
    out.save(str(out_path), deflate=True, garbage=3)
    out.close()
    if kind == "pdf":
        handle.close()

    print(f"[ok] {out_path}")
    print(f"[ok] lecture title: {lec_title!r}")
    print(f"[ok] slides = {n}, output PDF pages = {total_pages} "
          f"(≥{n}; extra pages = notes continuations)")
    return out_path, n, total_pages


def main(argv=None):
    ap = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("lecture_dir",
                    help="lecture folder, e.g. library/lectures/lec-04")
    ap.add_argument("--dpi", type=int, default=150,
                    help="raster DPI for slide images from the PDF (default 150)")
    ap.add_argument("--font", default=None,
                    help="override path to a Cyrillic-capable regular TTF")
    args = ap.parse_args(argv)

    lecture_dir = Path(args.lecture_dir).resolve()
    if not lecture_dir.exists():
        raise SystemExit(f"lecture dir not found: {lecture_dir}")
    build(lecture_dir, args.dpi, args.font)


if __name__ == "__main__":
    main()
