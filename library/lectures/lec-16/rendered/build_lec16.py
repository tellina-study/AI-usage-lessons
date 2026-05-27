"""
Full 43-slide build of Лекции 16 «AI в нефтегазовой отрасли и добыче ресурсов».

Source-of-truth: deck.yaml v1 + chapter v2.1 multi-part (~32k слов) + slides/*.md.

Issue #144 · downstream от chapter (book-first).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml ≈ 75 мин.

Lec-N-1 паттерн compliance: match lec-13/14 (cover + lecture-map + 7 section dividers +
keystone + выделенный Q&A; top progress bar только на dividers + cover).

Build via:
  python3 build_lec16.py            # generates lec-16.pptx (notes = file refs)
  python3 inject_notes.py           # injects FULL speaker notes from slides/*.md
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
RED_WARN = RGBColor(0xC0, 0x39, 0x2B)
ROADMAP = RGBColor(0xD9, 0xE2, 0xEC)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent
ASSETS = ROOT.parent / "assets"
OUT = ROOT / "lec-16.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *, size=16, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except: pass
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline_box(slide, x, y, w, h, lines, *, size=14, bold=False, color=DEEP,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
                  line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except: pass
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image_aspect(slide, path, x, y, w, h):
    """Add picture preserving aspect ratio (centered in box)."""
    p = Path(path)
    if not p.exists():
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        text_box(slide, x, y+h/2-0.3, w, 0.6, f"[нет: {p.name}]",
                 size=10, color=SLATE, align=PP_ALIGN.CENTER)
        return None
    try:
        with Image.open(p) as img:
            iw, ih = img.size
        img_ratio = iw / ih
        box_ratio = w / h
        if img_ratio > box_ratio:
            new_w = w
            new_h = w / img_ratio
            cx = x
            cy = y + (h - new_h) / 2
        else:
            new_h = h
            new_w = h * img_ratio
            cx = x + (w - new_w) / 2
            cy = y
        return slide.shapes.add_picture(str(p), Inches(cx), Inches(cy),
                                         width=Inches(new_w), height=Inches(new_h))
    except Exception as e:
        print(f"Image fail {path}: {e}")
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        return None


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# Roadmap sections for Lec-16: 7 разделов
SECTIONS = ["Стержень", "Q1 Зрелое", "Q3 Разведка", "Q2 Метан", "Q4 Новые опоры", "Россия", "Сквозные"]


def roadmap_bar(slide, current_section):
    """7-section roadmap bar at top of section dividers + cover."""
    bar_y = 0.20
    bar_h = 0.32
    total_w = 12.33
    seg_w = total_w / 7
    for i, name in enumerate(SECTIONS):
        x = 0.5 + i * seg_w
        is_active = (i == current_section)
        rectangle(slide, x, bar_y, seg_w - 0.05, bar_h,
                  fill=GOLD if is_active else ROADMAP)
        text_box(slide, x, bar_y + 0.04, seg_w - 0.05, bar_h - 0.04,
                 f"{i+1}. {name}", size=10, bold=is_active,
                 color=DEEP if is_active else SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, text, *, y=7.10):
    text_box(slide, 0.5, y, 12.33, 0.35, text,
             size=10, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)


def attribution(slide, text, *, x=0.5, y=6.95, w=12.33):
    text_box(slide, x, y, w, 0.3, text,
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)


def gold_callout(slide, x, y, w, h, text, *, size=14):
    rounded_box(slide, x, y, w, h, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, x+0.2, y+0.08, w-0.4, h-0.15, text,
             size=size, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


def section_divider(p, q_label, q_title, mood_line, tag_text, section_idx, large_size=160,
                    label_color=GOLD):
    """Section divider with large quadrant label + mood + tag."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=section_idx)
    # Decorative large Q-label (centered in upper area only)
    text_box(slide, 0.5, 1.0, 5.5, 3.6, q_label,
             size=large_size, bold=True, color=label_color, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    # Subtitle BELOW Q-label (clear separation)
    text_box(slide, 0.5, 4.85, 5.5, 0.55, q_title,
             size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    # Mood line
    multiline_box(slide, 6.2, 1.5, 6.6, 4.0, [
        (mood_line, {"size": 20, "bold": True, "color": DEEP}),
    ], line_spacing=1.3)
    # Bottom tag (gold tint)
    gold_callout(slide, 0.5, 6.0, 12.33, 0.85, tag_text, size=14)
    return slide


# ====================================================================
# SECTION 0: Введение + keystone (s01-s05)
# ====================================================================

def s01_hero_permian(p):
    """s01 — hero hook: Permian VIIRS satellite. Hero ≥40% area."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    img = ASSETS / "screenshots" / "s01-permian-viirs.jpg"
    add_image_aspect(slide, img, 0.5, 0.4, 7.5, 5.5)
    attribution(slide, "NASA Earth Observatory · VIIRS day-night band · NOAA · 2024",
                x=0.5, y=5.95, w=7.5)
    multiline_box(slide, 8.3, 0.6, 4.6, 5.4, [
        ("2 593", {"size": 72, "bold": True, "color": GOLD}),
        ("факельных шлейфа", {"size": 18, "bold": True, "color": DEEP}),
        ("Пермский бассейн, 2024", {"size": 14, "italic": True, "color": LIGHT}),
        ("", {"size": 10}),
        ("Пик: ~34 000 т/ч метана", {"size": 14, "bold": True, "color": DEEP}),
        ("из ~80 000 скважин (~3,2% факелуют)", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 10}),
        ("Не аварии — нормальный", {"size": 13, "color": DEEP}),
        ("эксплуатационный режим.", {"size": 13, "color": DEEP}),
        ("Сжигать избыточный газ", {"size": 13, "color": DEEP}),
        ("быстрее и дешевле, чем", {"size": 13, "color": DEEP}),
        ("строить газовую инфраструктуру.", {"size": 13, "color": DEEP}),
    ], line_spacing=1.1)
    gold_callout(slide, 0.5, 6.3, 12.33, 0.8,
                 "AI в нефтегазе — не «улучшалка на 5%», а способ закрыть конкретные провалы. И закрывает либо громко, либо публично проваливается.",
                 size=14)
    add_notes(slide, "См. slides/s01-hook-permian-viirs.md speaker notes.")


def s02_cover(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    roadmap_bar(slide, current_section=-1)
    text_box(slide, 0.5, 1.2, 4.0, 4.5, "16",
             size=240, bold=True, color=ROADMAP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    multiline_box(slide, 4.5, 1.5, 8.3, 2.6, [
        ("Лекция 16", {"size": 22, "bold": True, "color": LIGHT}),
        ("AI в нефтегазовой отрасли", {"size": 34, "bold": True, "color": DEEP}),
        ("и добыче ресурсов", {"size": 34, "bold": True, "color": DEEP}),
    ], line_spacing=1.05)
    text_box(slide, 4.5, 4.3, 8.3, 0.5,
             "Шесть разделов через матрицу данные × физика.",
             size=18, italic=True, color=TEAL)
    rounded_box(slide, 4.5, 4.95, 8.3, 1.85)
    multiline_box(slide, 4.7, 5.05, 8.0, 1.7, [
        ("Главный вопрос:", {"size": 12, "bold": True, "color": MID}),
        ("В каком квадранте матрицы «данные × физика» AI",
         {"size": 13, "color": DEEP}),
        ("становится мультипликатором классических методов,",
         {"size": 13, "color": DEEP}),
        ("а где он либо необходим, либо опасен — на материале",
         {"size": 13, "color": DEEP}),
        ("10 разобранных провалов индустрии 2014–2026?",
         {"size": 13, "bold": True, "color": DEEP}),
    ], line_spacing=1.15)
    footer(slide, "Курс «Применение AI в инженерии» · Студенты-инженеры 3 курса · 2026")
    add_notes(slide, "См. slides/s02-cover.md speaker notes.")


def s03_about(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.8,
             "Курс — про инженерное суждение «когда применять AI, когда отказаться»",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Сегодня — на материале нефтегазовой отрасли.",
             size=14, italic=True, color=LIGHT)
    cards = [
        ("Аудитория", "Студенты-инженеры\n3 курса\n(универсальная,\nне отраслевые\nспециалисты)", LIGHT),
        ("Формат", "42 слайда\n· 10 разобранных\nпровалов\n· 12+ рабочих\nкейсов\n· 7 разделов", MID),
        ("Вы научитесь", "Различать 4 квадранта\nматрицы\n· Критически читать\nвендорские заявления\n· Применять критерии\n«AI не нужен»\n· Альтернативы\nбез AI", TEAL),
    ]
    card_w = 4.0
    gap = 0.15
    x0 = 0.5
    y = 2.0
    for i, (title, body, accent) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        rounded_box(slide, x, y, card_w, 4.6)
        rectangle(slide, x, y, card_w, 0.7, fill=accent)
        text_box(slide, x+0.1, y, card_w-0.2, 0.7, title,
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x+0.2, y+0.9, card_w-0.4, 3.5, body,
                 size=14, color=DEEP, line_spacing=1.35)
    gold_callout(slide, 0.5, 6.85, 12.33, 0.5,
                 "10 разобранных провалов + 12+ рабочих кейсов = инженерный фильтр, не каталог инноваций.",
                 size=13)
    add_notes(slide, "См. slides/s03-about.md speaker notes.")


def s04_lecture_map(p):
    """s03 (new ID) — План лекции. Все 7 cards одной яркости (LIGHT) — это план, не «вы здесь»."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "План лекции — 7 разделов",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.05, 12.33, 0.4,
             "Не по цепочке создания стоимости, а по структуре AI-задачи: доступность данных × определённость процессов",
             size=14, italic=True, color=LIGHT)
    # Neutralize Q2 highlight — все 7 cards одной яркости (LIGHT).
    sections = [
        ("1", "Стержень", "Матрица", "Данные × процессы, 4 квадранта", LIGHT),
        ("2", "Q1 Зрелое", "3 кейса · 2 провала", "Зрелое производство", LIGHT),
        ("3", "Q3 Разведка", "3 кейса · 2 провала", "Разведка фронтиров", LIGHT),
        ("4", "Q2 Метан", "4 системы · 2 провала", "Метановая MRV", LIGHT),
        ("5", "Q4 Новые опоры", "2 пилота · 2 провала", "CCS + EGS", LIGHT),
        ("6", "Россия", "3 программы", "Санкции, инсорсинг", LIGHT),
        ("7", "Сквозные риски", "2 провала", "Кибер + кризис 2020", LIGHT),
    ]
    card_w = 1.71
    gap = 0.05
    x0 = 0.5
    y = 1.6
    for i, (num, title, dur, desc, accent) in enumerate(sections):
        x = x0 + i * (card_w + gap)
        rounded_box(slide, x, y, card_w, 4.4)
        circle(slide, x + card_w/2 - 0.28, y + 0.15, 0.55, 0.55, fill=accent)
        text_box(slide, x + card_w/2 - 0.28, y + 0.15, 0.55, 0.55, num,
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x + 0.1, y + 0.85, card_w - 0.2, 0.5, title,
                 size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x + 0.1, y + 1.35, card_w - 0.2, 0.4, dur,
                 size=9, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
        text_box(slide, x + 0.1, y + 1.85, card_w - 0.2, 2.5, desc,
                 size=11, color=DEEP, line_spacing=1.3, align=PP_ALIGN.CENTER)
    gold_callout(slide, 0.5, 6.2, 12.33, 0.85,
                 "Краткий словарь: MRV = выявление-учёт-проверка · OGI = оптическая газовая визуализация · CCS = улавливание и хранение углерода · EGS = улучшенные геотермальные системы · SIS = приборная система безопасности",
                 size=10)
    add_notes(slide, "См. slides/s03-lecture-map.md speaker notes.")


def s05_keystone_matrix(p):
    """s04 (new ID) — KEYSTONE: 2x2 matrix доступность данных × определённость процессов."""
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.3, 12.33, 0.7,
             "Когда AI работает в нефтегазе? Доступность данных × определённость процессов",
             size=22, bold=True, color=DEEP)
    text_box(slide, 0.5, 0.95, 12.33, 0.4,
             "От разведки фронтиров до спутникового метана — AI имеет 4 разных профиля",
             size=13, italic=True, color=LIGHT)
    # Y-axis labels (placed in column на левой стороне)
    text_box(slide, 0.1, 1.7, 1.5, 0.4, "Высокая ↑",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(slide, 0.05, 3.55, 1.6, 0.7, "Доступность\nданных",
             size=13, bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.1)
    text_box(slide, 0.1, 5.55, 1.5, 0.4, "↓ Низкая",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # X-axis labels
    text_box(slide, 2.0, 6.4, 3.5, 0.4, "← Низкая",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(slide, 5.0, 6.35, 3.5, 0.55, "Определённость\nпроцессов",
             size=13, bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.05)
    text_box(slide, 8.0, 6.4, 3.5, 0.4, "Высокая →",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # Matrix 2x2
    quad_x0 = 1.7
    quad_y0 = 1.55
    quad_w = 5.0
    quad_h = 2.25
    quads = [
        (0, 0, TEAL_TINT, "Q2 — Метановая MRV", "AI необходим\nMethaneSAT / Carbon Mapper / GHGSat", TEAL),
        (1, 0, SURFACE, "Q1 — Зрелое производство", "AI как мультипликатор\nAmbyint +15% на 200 скважинах", MID),
        (0, 1, GOLD_TINT, "Q4 — Новые опоры (CCS + EGS)", "AI и физика буксуют вместе\nNorthern Lights / Fervo", GOLD),
        (1, 1, SURFACE, "Q3 — Разведка фронтиров", "Сначала физика, AI как дополнение\nAramco METABRAIN / Eni HPC6", LIGHT),
    ]
    for col, row, fill, title, body, accent in quads:
        x = quad_x0 + col * (quad_w + 0.2)
        y = quad_y0 + row * (quad_h + 0.15)
        rounded_box(slide, x, y, quad_w, quad_h, fill=fill, stroke=accent, stroke_w=2.5)
        text_box(slide, x + 0.2, y + 0.1, quad_w - 0.4, 0.5, title,
                 size=15, bold=True, color=DEEP)
        text_box(slide, x + 0.2, y + 0.7, quad_w - 0.4, quad_h - 0.85, body,
                 size=12, color=DEEP, line_spacing=1.35)
    gold_callout(slide, 0.5, 7.0, 12.33, 0.45,
                 "За каждым AI-внедрением — альтернатива: физический симулятор, OGI-камера, классическая интерпретация.",
                 size=12)
    add_notes(slide, "См. slides/s04-keystone-matrix.md speaker notes.")


# ====================================================================
# SECTION 1: Q1 mainstream production (s06-s12)
# ====================================================================

def s06_q1_divider(p):
    return section_divider(
        p, "Q1", "Зрелое производство",
        "AI здесь — мультипликатор классических методов. Самый освоенный квадрант. И самый структурно проваленный.",
        "3 рабочих кейса · 2 структурных провала · 86% пилотов застряло — статистическая норма",
        section_idx=1, large_size=200, label_color=MID)


def s07_pilot_stuck(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.8,
             "86% AI-проектов в энергетике застряли в пилоте",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "McKinsey State of AI 2024 — vs средний по отраслям ~67%. Энергетика на 18 п.п. хуже среднего.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.7, 6.0, 4.6)
    img = ASSETS / "charts" / "s07-pilot-stuck.png"
    add_image_aspect(slide, img, 0.7, 1.85, 5.6, 4.3)
    rounded_box(slide, 6.7, 1.7, 6.13, 4.6)
    multiline_box(slide, 6.9, 1.85, 5.83, 4.4, [
        ("5 структурных причин:", {"size": 14, "bold": True, "color": MID}),
        ("", {"size": 4}),
        ("1. Данные. 60–80% времени AI-проекта = очистка данных. Только 21% энергокомпаний имеют качество для AI промышленного уровня.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("2. Интеграция со старой IT. Стоимость интеграции SCADA/MES/ERP — в 3–5× выше стоимости AI-софта.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("3. Дефицит кадров. AI + предметная область. После кризиса 2020 — 107 тыс. рабочих мест потеряно.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("4. Культура безопасности. Старший оператор отказывает рискованной рекомендации — и часто прав.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("5. Горизонт окупаемости. AI-поставщик обещает 12–18 мес; срок жизни месторождения 20–30 лет. Несовместимо.", {"size": 11, "color": DEEP}),
    ], line_spacing=1.2)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Это не «AI плохой» — это статистическая норма отрасли. 14% работают vs 86% застряли — инженерный фильтр, а не приговор.",
                 size=12)
    add_notes(slide, "См. slides/s06-86-percent-pilot-stuck.md speaker notes.")


def s07b_aspen_alert_fatigue(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.35, 12.33, 0.85,
             "«Усталость от ложных тревог устранена» — это маркетинг",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.2, 12.33, 0.4,
             "Aspen Mtell на нефтепереработке: 100–500 ложных тревог в день; общезаводские пилоты тихо закрываются.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.7, 5.5, 4.5)
    img = ASSETS / "screenshots" / "s09-aspen.jpg"
    add_image_aspect(slide, img, 0.7, 1.85, 5.1, 4.0)
    attribution(slide, "AspenTech · страница продукта Aspen Mtell", x=0.7, y=5.85, w=5.0)
    rounded_box(slide, 6.2, 1.7, 6.63, 4.5)
    multiline_box(slide, 6.4, 1.8, 6.3, 4.4, [
        ("Маркетинг (AspenTech):", {"size": 13, "bold": True, "color": MID}),
        ("«Снижение незапланированного простоя на 60%»", {"size": 11, "color": DEEP, "italic": True}),
        ("«Усталость операторов устранена»", {"size": 11, "color": DEEP, "italic": True}),
        ("", {"size": 4}),
        ("Реальность на НПЗ:", {"size": 13, "bold": True, "color": RED_WARN}),
        ("· 100–500 ложных тревог/день — оператор перестаёт реагировать", {"size": 11, "color": DEEP}),
        ("· Успех на одной колонне → общезаводской пилот тихо закрыт", {"size": 11, "color": DEEP}),
        ("· Honeywell UOP 310+ установок / ~700 НПЗ = ~44% охвата", {"size": 11, "color": DEEP}),
        ("· Многие НПЗ — «классический APC без AI»", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("Структурный разрыв:", {"size": 13, "bold": True, "color": GOLD}),
        ("Многослойная физика (масса + энергия + реакция + коррозия) ломает ML-суррогаты на нестандартных режимах.", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    gold_callout(slide, 0.5, 6.35, 12.33, 0.55,
                 "Урок: отчёт поставщика ≠ реальность на установке. Коммерческий учёт, SIS, общезаводская переработка — AI ещё не дошёл.",
                 size=12)
    add_notes(slide, "См. slides/s07-aspen-alert-fatigue.md speaker notes.")


def s08_ambyint(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "Ambyint InfinityRL: +15% на 200 скважинах Permian",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Канадский стартап Калгари (2014), $25 млн раунд Series B 2022. Обучение с подкреплением для оптимизации механизированной добычи.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.8, 6.0, 4.5)
    img = ASSETS / "charts" / "s08-ambyint-delta.png"
    add_image_aspect(slide, img, 0.7, 1.95, 5.6, 4.2)
    rounded_box(slide, 6.7, 1.8, 6.13, 4.5)
    multiline_box(slide, 6.9, 1.95, 5.83, 4.3, [
        ("Метрики кейса:", {"size": 14, "bold": True, "color": MID}),
        ("· Регионы: Permian + Eagle Ford + Bakken", {"size": 12, "color": DEEP}),
        ("· Тип: штанговые насосы + ЭЦН (зрелая добыча)", {"size": 12, "color": DEEP}),
        ("· Исходный уровень: 100–500 баррелей/день на скважину", {"size": 12, "color": DEEP}),
        ("· Прирост: +15% от среднего исторического", {"size": 12, "color": DEEP, "bold": True}),
        ("· На 200 скважинах: +3 000–15 000 баррелей/день суммарно", {"size": 12, "color": GOLD, "bold": True}),
        ("", {"size": 6}),
        ("Почему это сильный кейс:", {"size": 14, "bold": True, "color": MID}),
        ("1. Проверяемая база — не «спас $10 млн» без знаменателя.", {"size": 12, "color": DEEP}),
        ("2. Обучение с подкреплением поверх классики — дополнение, не замена.", {"size": 12, "color": DEEP}),
        ("3. Узкая область — оптимизация подъёма, не «AI везде».", {"size": 12, "color": DEEP}),
    ], line_spacing=1.3)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Когда Ambyint НЕ работает: стрипперные скважины (<10 баррелей/день). +15% = +1,5 барр./день; стоимость развёртывания > извлечённой ценности.",
                 size=12)
    add_notes(slide, "См. slides/s08-ambyint-infinityrl.md speaker notes.")


def s09_vendor_landscape(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Ландшафт поставщиков Q1 — 3 группы",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.05, 12.33, 0.4,
             "Тот же режим (прогностическое обслуживание) реализуют 10+ поставщиков. Режим ≠ бренд.",
             size=13, italic=True, color=LIGHT)
    groups = [
        ("ML-стартапы\nдля добычи", LIGHT, [
            "Ambyint (Калгари) — InfinityRL для\nмеханизированной добычи",
            "OspreyData (→ Mesquite 2022) —\nML с экспертом для ЭЦН, газлифта",
        ]),
        ("Корпоративные\n(нацкомпании + супермэйджоры)", MID, [
            "SLB Avocet — ПО управления\nдобычей с ML с 2020+",
            "Halliburton DecisionSpace\nProduction — аналог. Промышленный масштаб.",
        ]),
        ("НПЗ + трубопроводы", TEAL, [
            "AspenTech Aspen Mtell (Emerson)\n— растёт через перекрёстные продажи",
            "Honeywell UOP Connect — 310+\nустановок / ~700 НПЗ = ~44%",
            "Yokogawa / ABB Genix / Emerson —\nконкуренты",
        ]),
    ]
    col_w = 4.0
    gap = 0.15
    x0 = 0.5
    y = 1.6
    for i, (title, accent, items) in enumerate(groups):
        x = x0 + i * (col_w + gap)
        rounded_box(slide, x, y, col_w, 4.3)
        rectangle(slide, x, y, col_w, 0.95, fill=accent)
        text_box(slide, x+0.1, y+0.1, col_w-0.2, 0.8, title,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        lines = []
        for j, item in enumerate(items):
            lines.append((f"· {item}", {"size": 11, "color": DEEP}))
            if j < len(items) - 1:
                lines.append(("", {"size": 4}))
        multiline_box(slide, x+0.2, y+1.1, col_w-0.4, 3.0, lines, line_spacing=1.3)
    gold_callout(slide, 0.5, 6.05, 12.33, 0.85,
                 "Бурение (дополнительно): Nabors PACE-X (4-мильный ствол Bakken) · NOV NOVOS · Precision Drilling AlphaAutomation. Режим = автоматизация бурения; бренд вторичен.",
                 size=12)
    add_notes(slide, "См. slides/s09-q1-vendor-landscape.md speaker notes.")


def s10_rosneft_digital_field(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "Роснефть Digital Field на Башнефть Илишевское",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.25, 12.33, 0.4,
             "Внутренняя разработка после ухода Roxar / Schlumberger в 2022. Вертикальная интеграция по умолчанию.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.8, 6.0, 4.5)
    multiline_box(slide, 0.7, 1.95, 5.7, 4.3, [
        ("Метрики Илишевского, 2024:", {"size": 14, "bold": True, "color": MID}),
        ("", {"size": 6}),
        ("+1 Mt/год", {"size": 36, "bold": True, "color": GOLD}),
        ("дополнительной нефти (= +5,9% от ~17 Mt/год)", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("~1 млрд ₽/год", {"size": 24, "bold": True, "color": DEEP}),
        ("дополнительного эффекта", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("23 продукта (10 коммерциализованных)", {"size": 14, "color": DEEP}),
        ("+60% удалённого управления · +5% энергоэфф. · −5% логистики", {"size": 11, "color": DEEP}),
    ], line_spacing=1.2)
    rounded_box(slide, 6.7, 1.8, 6.13, 4.5)
    multiline_box(slide, 6.9, 1.95, 5.83, 4.3, [
        ("Контекст: санкции → внутренняя разработка", {"size": 14, "bold": True, "color": MID}),
        ("", {"size": 6}),
        ("· После марта 2022 — Roxar (Schlumberger), AspenTech, Honeywell ушли с РФ.", {"size": 12, "color": DEEP}),
        ("", {"size": 4}),
        ("· Вертикальная интеграция — необходимость, не выбор.", {"size": 12, "color": DEEP}),
        ("", {"size": 4}),
        ("· Ближе к китайской модели (Sinopec, CNOOC), чем к американской на основе поставщиков.", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Оговорка:", {"size": 12, "bold": True, "color": GOLD}),
        ("Российские KPI — самоотчёт в пресс-релизе; независимый аудит ограничен санкциями. Тот же уровень осторожности, что для Aramco $1,8 млрд.", {"size": 11, "italic": True, "color": DEEP}),
    ], line_spacing=1.3)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Башнефть Илишевское — типичное зрелое месторождение Q1: эксплуатируется с 1980-х, ~17 млн т/год исходный уровень. Пилотировать AI на известном активе, не на фронтире.",
                 size=12)
    add_notes(slide, "См. slides/s10-rosneft-digital-field.md speaker notes.")


def s11_cognite_c3ai(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "Cognite + C3.ai: чистые AI-вендоры теряют долю",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.25, 12.33, 0.4,
             "Базовые модели поедают узких отраслевых AI-вендоров. Адресуемый рынок уже, чем ожидалось.",
             size=13, italic=True, color=LIGHT)
    boxes = [
        ("Cognite (Норвегия, выделён из Aker BP 2017)", LIGHT, [
            "2018: оценка ~$300 млн",
            "2021–2022: план IPO $2–3 млрд в 2023",
            "2023: IPO отменён — рыночные условия + сжигание капитала",
            "2024: ARR (годовая повторяющаяся выручка) $94 млн (+40%), 871 сотрудник после реструктуризации",
            "2026: «время IPO неопределённо» (отчёт Aker ASA)",
        ]),
        ("C3.ai (США, основан 2009 Tom Siebel)", TEAL, [
            "BHC3 JV с Baker Hughes (2019) — реструктурирован к 2023",
            "ФГ24 нефтегаз = 5,9% выручки = ~$18 млн из $310 млн",
            "ФГ25: «не-нефтегазовая выручка +48%» → нефтегаз уменьшается",
            "C3.ai сместил фокус на федеральный/оборонный сектор",
            "IPO 2020: $42 → пик ~$15 млрд → 2026: ~$3–4 млрд (−70–80% от пика)",
        ]),
    ]
    col_w = 6.0
    gap = 0.33
    x0 = 0.5
    y = 1.8
    for i, (title, accent, items) in enumerate(boxes):
        x = x0 + i * (col_w + gap)
        rounded_box(slide, x, y, col_w, 4.5)
        rectangle(slide, x, y, col_w, 0.7, fill=accent)
        text_box(slide, x+0.15, y, col_w-0.3, 0.7, title,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
        lines = []
        for item in items:
            lines.append((f"· {item}", {"size": 12, "color": DEEP}))
            lines.append(("", {"size": 4}))
        multiline_box(slide, x+0.2, y+0.85, col_w-0.4, 3.5, lines, line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Базовые модели (SLB Lumi, METABRAIN) поедают узкоспециализированных: (a) большая модель + лёгкое дообучение (LoRA) > специализированная; (b) крупные нацкомпании разрабатывают свои базовые модели.",
                 size=12)
    add_notes(slide, "См. slides/s11-cognite-c3ai-decline.md speaker notes.")


def s12_q1_no_ai_criteria(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.8,
             "Когда AI НЕ нужен в Q1 — 6 структурных критериев",
             size=24, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Дистиллировано из практики последних 5 лет. Каждый — с конкретной альтернативой.",
             size=13, italic=True, color=LIGHT)
    criteria = [
        ("1", "Зрелый пласт + классический симулятор (Eclipse)", "Опытный инженер + классический симулятор дают надёжные ответы. ML — накладные расходы без существенного прироста.", LIGHT),
        ("2", "Стрипперные скважины (<10 барр./день)", "Прибавка +15% = +1,5 барр./день; стоимость развёртывания > извлечённой ценности. Юнит-экономика отрицательная.", GOLD),
        ("3", "Коммерческий учёт нефти (custody transfer)", "Передача товарной нефти. Регулятор требует расходомер 0,2% точности. Не «чёрный ящик» ML.", MID),
        ("4", "Аварийная остановка (BOP/PRV/ESD) — SIS", "SIL3/SIL4 по IEC 61511 = детерминировано + сертифицируемо. ML не сертифицируется.", TEAL),
        ("5", "Разведка фронтиров без аналогов", "ML не на чем обучать. Опытный геофизик + классическая интерпретация (далее в Разделе 2).", DEEP),
        ("6", "Отчётность EU Methane Reg", "Прослеживаемость обязательна — не оценка «чёрного ящика». Регуляторное соответствие.", LIGHT),
    ]
    card_w = 4.0
    card_h = 2.3
    gap_x = 0.15
    gap_y = 0.15
    x0 = 0.5
    y0 = 1.6
    for i, (num, title, body, accent) in enumerate(criteria):
        col = i % 3
        row = i // 3
        x = x0 + col * (card_w + gap_x)
        y = y0 + row * (card_h + gap_y)
        rounded_box(slide, x, y, card_w, card_h, stroke=accent, stroke_w=2)
        circle(slide, x + 0.15, y + 0.15, 0.6, 0.6, fill=accent)
        text_box(slide, x + 0.15, y + 0.15, 0.6, 0.6, num,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x + 0.85, y + 0.2, card_w - 1.0, 0.5, title,
                 size=13, bold=True, color=DEEP)
        text_box(slide, x + 0.2, y + 0.85, card_w - 0.4, card_h - 1.0, body,
                 size=10, color=DEEP, line_spacing=1.3)
    gold_callout(slide, 0.5, 6.65, 12.33, 0.45,
                 "Главный навык курса: уметь сказать «нет» там, где AI не нужен. 14% работают vs 86% застряли — разница часто именно здесь.",
                 size=12)
    add_notes(slide, "См. slides/s12-q1-no-ai-criteria.md speaker notes.")


# ====================================================================
# SECTION 2: Q3 frontier exploration (s13-s19)
# ====================================================================

def s13_q3_divider(p):
    return section_divider(
        p, "Q3", "Разведка фронтиров",
        "Каждая поисковая скважина = $50–100 млн. Размер выборки 1–5 скважин. ML не обобщается без аналогов. Физика — эталон, AI как дополнение.",
        "3 рабочих кейса · 2 провала десятилетия · гонка суперкомпьютеров $100–400 млн на инсталляцию",
        section_idx=2, large_size=200, label_color=LIGHT)


def s14_hpc_eni_aramco(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "Гонка суперкомпьютеров (HPC) Q3: $100–400 млн на инсталляцию",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Не превращается в товар как облако — стратегические капитальные затраты. Малые операторы вытесняются капитальным барьером.",
             size=13, italic=True, color=LIGHT)
    # Eni HPC6
    rounded_box(slide, 0.5, 1.85, 6.0, 4.5)
    rectangle(slide, 0.5, 1.85, 6.0, 0.55, fill=MID)
    text_box(slide, 0.65, 1.85, 5.7, 0.55, "Eni HPC6 (декабрь 2024)",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    img = ASSETS / "screenshots" / "s14-eni.jpg"
    add_image_aspect(slide, img, 0.65, 2.5, 5.7, 1.7)
    multiline_box(slide, 0.65, 4.3, 5.7, 2.05, [
        ("606 PFLOPS пиковая · 477 устойчивая", {"size": 13, "color": DEEP, "bold": True}),
        ("14 000 AMD MI250X GPU", {"size": 12, "color": DEEP}),
        ("$104 млн капитальных затрат", {"size": 14, "bold": True, "color": GOLD}),
        ("Top500 #5 мирового рейтинга", {"size": 12, "color": DEEP}),
        ("Стратегия: дешевле за FLOP (AMD vs NVIDIA)", {"size": 11, "color": SLATE, "italic": True}),
    ], line_spacing=1.3)
    # Aramco METABRAIN
    rounded_box(slide, 6.7, 1.85, 6.13, 4.5)
    rectangle(slide, 6.7, 1.85, 6.13, 0.55, fill=TEAL)
    text_box(slide, 6.85, 1.85, 5.83, 0.55, "Aramco METABRAIN (2024–2025)",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    img2 = ASSETS / "charts" / "s14-aramco-roi.png"
    add_image_aspect(slide, img2, 6.85, 2.5, 5.83, 1.7)
    multiline_box(slide, 6.85, 4.3, 5.83, 2.05, [
        ("~250 млрд параметров", {"size": 13, "color": DEEP, "bold": True}),
        ("7 трлн токенов = 90 лет данных Aramco", {"size": 12, "color": DEEP}),
        ("$1,8 млрд реализовано в 2024 (Давос янв 2025)", {"size": 14, "bold": True, "color": GOLD}),
        ("6 000 сотрудников обучены, 430 сценариев применения", {"size": 12, "color": DEEP}),
        ("Источник: CEO Amin Nasser", {"size": 11, "color": SLATE, "italic": True}),
    ], line_spacing=1.3)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Aramco выручка 2024 = $436,6 млрд → $1,8 млрд / $436,6 млрд = 0,41%. AI добавляет полпроцента к полностью оптимизированной операции.",
                 size=12)
    add_notes(slide, "См. slides/s14-hpc-eni-aramco.md speaker notes.")


def s15_slb_lumi(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "SLB Lumi (сентябрь 2024) — базовая модель поверх Petrel + Delfi",
             size=24, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Якорные заказчики: Aker BP, Shell, Azule Energy. Вычислительная база — NVIDIA Grace Hopper.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.85, 5.5, 4.4)
    img = ASSETS / "screenshots" / "s15-slb2.jpg"
    add_image_aspect(slide, img, 0.65, 2.0, 5.2, 4.05)
    attribution(slide, "SLB / Schlumberger пресс-кит · 2024", x=0.65, y=6.05, w=5.2)
    rounded_box(slide, 6.2, 1.85, 6.63, 4.4)
    multiline_box(slide, 6.4, 1.95, 6.3, 4.25, [
        ("Что Lumi делает:", {"size": 14, "bold": True, "color": MID}),
        ("· Интерпретация петрофизики — каротаж", {"size": 12, "color": DEEP}),
        ("· Автотрассировка сейсмических горизонтов", {"size": 12, "color": DEEP}),
        ("· Предварительная характеристика коллектора — фации", {"size": 12, "color": DEEP}),
        ("· Оптимизация параметров бурения", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Контекст:", {"size": 14, "bold": True, "color": MID}),
        ("· Цифровая выручка SLB 2024 = $2+ млрд = 5,7% от $35 млрд", {"size": 12, "color": DEEP}),
        ("· Halliburton + OpenAI/Anthropic — через партнёрства", {"size": 12, "color": DEEP}),
        ("· Режим «отраслевая базовая модель» ≠ бренд «Lumi»", {"size": 12, "color": GOLD, "bold": True}),
    ], line_spacing=1.3)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Через 2-3 года: SLB Lumi 2.0 или конкурент. Бренд вторичен, режим (базовая модель на 80+ лет отраслевых данных) первичен.",
                 size=12)
    add_notes(slide, "См. slides/s15-slb-lumi.md speaker notes.")


def s16_exxon_discovery6(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "ExxonMobil Discovery 6 — 4D-сейсмика месяцы → недели",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "HPE Cray EX4000, 4 032 NVIDIA Grace Hopper, 4× вычислений vs Discovery 5. $200–400 млн капитальных затрат.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.85, 5.5, 4.4)
    img = ASSETS / "screenshots" / "s16-stabroek.jpg"
    add_image_aspect(slide, img, 0.65, 2.0, 5.2, 4.05)
    attribution(slide, "USS Normandy у FPSO в Stabroek Block · US Navy / Wikimedia Commons (PD)", x=0.65, y=6.05, w=5.2)
    rounded_box(slide, 6.2, 1.85, 6.63, 4.4)
    multiline_box(slide, 6.4, 1.95, 6.3, 4.25, [
        ("Что делает Discovery 6:", {"size": 14, "bold": True, "color": MID}),
        ("· 4D-сейсмика (3D + время) — месяцы → недели", {"size": 12, "color": DEEP, "bold": True}),
        ("· Активное управление пластом в реальном времени", {"size": 12, "color": DEEP}),
        ("· Stabroek Block Guyana — основной сценарий", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Кейс Stabroek:", {"size": 14, "bold": True, "color": MID}),
        ("· 6 плавучих платформ (FPSO) к 2026", {"size": 12, "color": DEEP}),
        ("$1 млрд+", {"size": 28, "bold": True, "color": GOLD}),
        ("высвобождено через быстрое перепозиционирование скважин", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 6}),
        ("Параллель Aramco METABRAIN: тот же режим, разные стратегии (США vs Саудовская Аравия).", {"size": 11, "color": DEEP, "italic": True}),
    ], line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Капитальный барьер: $200–400 млн для топового суперкомпьютера исключает 95% операторов. Только нацкомпании + супермэйджоры.",
                 size=12)
    add_notes(slide, "См. slides/s16-exxonmobil-discovery6.md speaker notes.")


def s17_bp_beyond_limits(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "BP + Beyond Limits — $20 млн, разворот поставщика в 2023",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "BP Ventures, раунд Series B июнь 2017 → Beyond Limits разворот в медицину/производство 2023.",
             size=13, italic=True, color=LIGHT)
    # Left: promises
    rounded_box(slide, 0.5, 1.85, 6.0, 4.5)
    rectangle(slide, 0.5, 1.85, 6.0, 0.55, fill=LIGHT)
    text_box(slide, 0.65, 1.85, 5.7, 0.55, "Что обещали (2018):",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 2.55, 5.6, 3.7, [
        ("«AI впитает знания геологов", {"size": 13, "italic": True, "color": DEEP}),
        ("и будет имитировать их решения».", {"size": 13, "italic": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Происхождение:", {"size": 13, "bold": True, "color": MID}),
        ("· NASA JPL (наследие $20 млн)", {"size": 12, "color": DEEP}),
        ("· Glendale, Калифорния, 2014", {"size": 12, "color": DEEP}),
        ("· «Когнитивный AI» как отличие", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Сигналы BP 2018:", {"size": 13, "bold": True, "color": MID}),
        ("· $20 млн Series B июнь 2017", {"size": 12, "color": DEEP}),
        ("· Попытка ребрендинга Beyond Petroleum", {"size": 12, "color": DEEP}),
        ("· Публичное обязательство", {"size": 12, "color": DEEP}),
    ], line_spacing=1.3)
    # Right: failures
    rounded_box(slide, 6.7, 1.85, 6.13, 4.5)
    rectangle(slide, 6.7, 1.85, 6.13, 0.55, fill=RED_WARN)
    text_box(slide, 6.85, 1.85, 5.83, 0.55, "Что получили (2018–2023):",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 6.9, 2.55, 5.73, 3.7, [
        ("7 лет публичных", {"size": 36, "bold": True, "color": GOLD}),
        ("результатов нет", {"size": 18, "bold": True, "color": DEEP}),
        ("", {"size": 6}),
        ("· 0 кейсов на сайте BP после 2019", {"size": 12, "color": DEEP}),
        ("· 0 публикаций в Society of Petroleum Engineers", {"size": 12, "color": DEEP}),
        ("· Beyond Limits разворот в медицину/производство 2023", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("3 урока:", {"size": 13, "bold": True, "color": RED_WARN}),
        ("1. Концентрация на одном малом поставщике — риск", {"size": 11, "color": DEEP}),
        ("2. Когнитивный маркетинг — антропоморфное обещание", {"size": 11, "color": DEEP}),
        ("3. Рамка имитации — AI не имитирует, он аппроксимирует", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "«Когнитивный AI имитирует геолога» — антропоморфная рамка, скрывающая оценку структурной применимости.",
                 size=12)
    add_notes(slide, "См. slides/s17-bp-beyond-limits-failure.md speaker notes.")


def s18_ibm_repsol(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "IBM Watson + Repsol Kalimba (2014–2022)",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Цикл хайпа 2014–2016 → реальное коммерческое применение ≤10% от ожиданий. Watson Industry Solutions стагнация 2018–2022.",
             size=13, italic=True, color=LIGHT)
    # Kalimba project
    rounded_box(slide, 0.5, 1.85, 6.0, 4.5)
    rectangle(slide, 0.5, 1.85, 6.0, 0.55, fill=LIGHT)
    text_box(slide, 0.65, 1.85, 5.7, 0.55, "Проект Kalimba 2014–2022",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 2.55, 5.6, 3.7, [
        ("2014:", {"size": 13, "bold": True, "color": MID}),
        ("· IBM + Repsol объявляют партнёрство", {"size": 12, "color": DEEP}),
        ("· Cognitive Environments Lab (Нью-Йорк) + Repsol Tech Centre (Мадрид)", {"size": 12, "color": DEEP}),
        ("· «Когнитивные вычисления для разведки»", {"size": 12, "color": DEEP, "italic": True}),
        ("", {"size": 6}),
        ("2014–2022:", {"size": 13, "bold": True, "color": MID}),
        ("· 30 лет данных разведки «проанализированы»", {"size": 12, "color": DEEP}),
        ("· Конкретных новых открытий не объявлено", {"size": 12, "color": DEEP}),
        ("· Бюджет разведки Repsol сократился с 2019", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("2022:", {"size": 13, "bold": True, "color": RED_WARN}),
        ("· Тихое сворачивание", {"size": 12, "color": DEEP, "bold": True}),
    ], line_spacing=1.3)
    # Watson Health parallel
    rounded_box(slide, 6.7, 1.85, 6.13, 4.5)
    rectangle(slide, 6.7, 1.85, 6.13, 0.55, fill=RED_WARN)
    text_box(slide, 6.85, 1.85, 5.83, 0.55, "Watson Health parallel",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 6.9, 2.55, 5.73, 3.7, [
        ("$5 млрд+", {"size": 32, "bold": True, "color": GOLD}),
        ("инвестиций IBM в Watson Health 2015–2021", {"size": 12, "color": DEEP, "italic": True}),
        ("", {"size": 8}),
        ("→ продан Francisco Partners в 2022", {"size": 14, "bold": True, "color": DEEP}),
        ("за ~$1 млрд (= 20% от инвестиций)", {"size": 12, "color": DEEP}),
        ("", {"size": 8}),
        ("3 урока:", {"size": 13, "bold": True, "color": RED_WARN}),
        ("1. Универсальный «когнитивный» AI не масштабируется в узкую область", {"size": 11, "color": DEEP}),
        ("2. Цикл хайпа 2014–2016 → реальное применение ≤10% от ожиданий", {"size": 11, "color": DEEP}),
        ("3. SLB Lumi и Aramco METABRAIN — отраслевые = успешнее", {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Watson Health → ~$1 млрд продан = $5 млрд инвестиций vs $1 млрд продажа = провал $4 млрд. Косвенный якорь: Kalimba не был исключением — он был частью паттерна.",
                 size=12)
    add_notes(slide, "См. slides/s18-ibm-repsol-failure.md speaker notes.")


def s19_q3_alternatives(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "Альтернатива Q3: физические симуляторы + старшая экспертиза",
             size=24, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "AI как дополнение поверх, не как замена. Через 5 лет Eclipse будет стандартом — не AI-замена.",
             size=13, italic=True, color=LIGHT)
    sims = [
        ("Eclipse (SLB)", "Отраслевой стандарт симулятора пласта с 1983.\nЗрелые коллекторы + регуляторные отчёты.\nСвязанные уравнения флюид + теплопередача + химия.", MID),
        ("INTERSECT (SLB)", "Преемник для массивно-параллельных суперкомпьютеров.\nКрупные модели, более мелкая сетка.\nЗаменяет Eclipse на топовых проектах.", LIGHT),
        ("CMG IMEX / STARS / GEM", "Computer Modelling Group (Калгари).\nIMEX «чёрная нефть», STARS тепловые/МУН, GEM композиционный.\nЛидер в тяжёлой нефти + МУН (методы увеличения нефтеотдачи).", TEAL),
        ("OpenFOAM (CFD)", "Открытый CFD для расчётов NPV,\nгеомеханики, моделирования ГРП.\nИсследования + характеристика коллектора.", DEEP),
    ]
    sim_w = 6.0
    sim_h = 2.0
    gap = 0.2
    x0 = 0.5
    y0 = 1.85
    for i, (name, body, accent) in enumerate(sims):
        col = i % 2
        row = i // 2
        x = x0 + col * (sim_w + gap)
        y = y0 + row * (sim_h + gap)
        rounded_box(slide, x, y, sim_w, sim_h, stroke=accent, stroke_w=2)
        rectangle(slide, x, y, 0.15, sim_h, fill=accent)
        text_box(slide, x + 0.3, y + 0.1, sim_w - 0.4, 0.5, name,
                 size=14, bold=True, color=DEEP)
        text_box(slide, x + 0.3, y + 0.6, sim_w - 0.4, sim_h - 0.7, body,
                 size=11, color=DEEP, line_spacing=1.35)
    gold_callout(slide, 0.5, 6.15, 12.33, 0.85,
                 "Старший геофизик ($200–500 тыс./год) + классическая интерпретация > $5–20 млн базовая модель на фронтире. PINN (нейросеть с встроенной физикой) — исследовательский уровень, не коммерческий. AI как дополнение, не замена.",
                 size=12)
    add_notes(slide, "См. slides/s19-q3-alternatives.md speaker notes.")


# ====================================================================
# SECTION 3: Q2 methane MRV (s20-s27)
# ====================================================================

def s20_methane_alphabet(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Алфавит метановой MRV — 6 обязательных терминов",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.05, 12.33, 0.4,
             "Все встретятся в следующих 8 слайдах. Названия брендов — не переводим.",
             size=13, italic=True, color=LIGHT)
    terms = [
        ("MRV", "Monitoring / Reporting / Verification\n= выявление-учёт-проверка\nОбязательно под EU 2024/1787.", TEAL),
        ("OGI", "Optical Gas Imaging\n= оптическая газовая визуализация\nИК-камера (FLIR GFx320, Opgal EyeCGas).", MID),
        ("LDAR", "Leak Detection And Repair\n= программа выявления и устранения утечек\n4×/год обходы с OGI + ремонт 5-15 дней.", LIGHT),
        ("OGMP 2.0", "Oil & Gas Methane Partnership 2.0 (UN)\nУровень 5 — прямое измерение.\n170+ компаний подписали.", GOLD),
        ("SIL / SIS", "Safety Integrity Level / Safety Instrumented System\nIEC 61511. SIL3 = 0,001–0,0001 PFD.\nML не сертифицируется.", DEEP),
        ("барр./день / ЭЦН", "Баррелей нефти в день / Электроцентробежный насос\n(Electric Submersible Pump)\nСтрипперные скважины <10 барр./день.", LIGHT),
    ]
    card_w = 4.0
    card_h = 2.3
    gap = 0.15
    x0 = 0.5
    y0 = 1.6
    for i, (acro, body, accent) in enumerate(terms):
        col = i % 3
        row = i // 3
        x = x0 + col * (card_w + gap)
        y = y0 + row * (card_h + gap)
        rounded_box(slide, x, y, card_w, card_h, stroke=accent, stroke_w=2)
        rectangle(slide, x, y, card_w, 0.6, fill=accent)
        text_box(slide, x+0.15, y+0.05, card_w-0.3, 0.5, acro,
                 size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x+0.2, y+0.7, card_w-0.4, card_h - 0.85, body,
                 size=11, color=DEEP, line_spacing=1.3)
    gold_callout(slide, 0.5, 6.65, 12.33, 0.45,
                 "Дополнительно: AI необходим именно потому, что OGI + Picarro + LI-COR покрывают уровень площадки — точечное выявление утечки требует слияния 4 сенсоров.",
                 size=12)
    add_notes(slide, "См. slides/s20-methane-alphabet.md speaker notes.")


def s21_q2_divider(p):
    return section_divider(
        p, "Q2", "Метановая MRV",
        "Данные — петабайты в день. Физика — разорвана. Слияние 4 сенсоров + атрибуция малой утечки — открытая ML-задача. AI необходим. Но один спутник = катастрофическая единичная уязвимость.",
        "4 рабочих системы · 2 провала · регуляторное давление со стороны EU 2024/1787",
        section_idx=3, large_size=200, label_color=TEAL)


def s22_methanesat_permian(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "MethaneSAT (4 марта 2024) — ключевой результат по Permian",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Первый в истории спутник, владелец которого — экологическая некоммерческая организация. Бюджет $88 млн.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.85, 6.0, 4.5)
    img = ASSETS / "screenshots" / "s22-methanesat.png"
    add_image_aspect(slide, img, 0.65, 2.0, 5.7, 3.5)
    attribution(slide, "MethaneSAT / EDF · 2024", x=0.65, y=5.55, w=5.7)
    multiline_box(slide, 0.65, 5.85, 5.7, 0.45, [
        ("Что отличало MethaneSAT:", {"size": 12, "bold": True, "color": MID}),
    ])
    rounded_box(slide, 6.7, 1.85, 6.13, 4.5)
    multiline_box(slide, 6.9, 1.95, 5.83, 4.3, [
        ("Технические возможности:", {"size": 13, "bold": True, "color": MID}),
        ("· Широкий охват — 200×200 км за один проход", {"size": 11, "color": DEEP}),
        ("· Высокая точность — порог детекции ~500 кг/ч", {"size": 11, "color": DEEP}),
        ("· Открытый доступ через Google Earth Engine", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Ключевой результат по Permian:", {"size": 13, "bold": True, "color": MID}),
        ("410 т/ч", {"size": 32, "bold": True, "color": GOLD}),
        ("метана = +50% над оценкой EPA (~273 т/ч)", {"size": 12, "italic": True, "color": DEEP}),
        ("", {"size": 4}),
        ("· Нью-Мексико: 1,2% интенсивности утечек", {"size": 11, "color": DEEP}),
        ("· Техас: 3,1% интенсивности утечек", {"size": 11, "color": DEEP}),
        ("· ~2 000 файлов данных за 15,5 мес работы", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "+50% над EPA — это не «AI ошибается». Это AI измеряет, а EPA считает по коэффициентам эмиссий 30-летней давности.",
                 size=12)
    add_notes(slide, "См. slides/s22-methanesat-permian.md speaker notes.")


def s23_methanesat_loss(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "20 июня 2025 — MethaneSAT потерян после 15,5 месяцев",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "4 марта 2024 запуск → 20 июня 2025 «аномалия аппарата». 26% от расчётного 5-летнего срока.",
             size=13, italic=True, color=LIGHT)
    # Timeline LEFT
    rounded_box(slide, 0.5, 1.85, 6.0, 4.5)
    rectangle(slide, 0.5, 1.85, 6.0, 0.55, fill=LIGHT)
    text_box(slide, 0.65, 1.85, 5.7, 0.55, "Хронология",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 2.55, 5.6, 3.75, [
        ("4 марта 2024", {"size": 14, "bold": True, "color": MID}),
        ("Запуск SpaceX Falcon 9", {"size": 12, "color": DEEP}),
        ("", {"size": 4}),
        ("Март 2024 — июнь 2025", {"size": 14, "bold": True, "color": MID}),
        ("15,5 месяцев операций", {"size": 12, "color": DEEP}),
        ("~2 000 файлов данных собрано", {"size": 12, "color": DEEP}),
        ("", {"size": 4}),
        ("20 июня 2025", {"size": 14, "bold": True, "color": RED_WARN}),
        ("«Аномалия аппарата» — потеря связи", {"size": 12, "color": DEEP}),
        ("", {"size": 4}),
        ("15,5 / 60 мес расчётных", {"size": 14, "bold": True, "color": GOLD}),
        ("= 26% реализованного срока", {"size": 12, "color": DEEP}),
        ("", {"size": 4}),
        ("$5,7 млн/мес фактически vs $1,5 млн/мес планировалось", {"size": 11, "color": DEEP, "italic": True}),
    ], line_spacing=1.2)
    # Lessons RIGHT
    rounded_box(slide, 6.7, 1.85, 6.13, 4.5)
    rectangle(slide, 6.7, 1.85, 6.13, 0.55, fill=RED_WARN)
    text_box(slide, 6.85, 1.85, 5.83, 0.55, "4 урока потери",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 6.9, 2.55, 5.83, 3.75, [
        ("1. Единая точка отказа (SPOF)", {"size": 13, "bold": True, "color": MID}),
        ("Один спутник = катастрофический вектор потери. Группировка (GHGSat 13+) — устойчива.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("2. Надёжность аппаратуры", {"size": 13, "bold": True, "color": MID}),
        ("Открытый космос ≠ контролируемая лаборатория. 26% срока — высокий риск.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("3. Регуляторное ограничение", {"size": 13, "bold": True, "color": MID}),
        ("Инспектор по EU 2024/1787 не может полагаться на один источник. Диверсификация обязательна.", {"size": 11, "color": DEEP}),
        ("", {"size": 4}),
        ("4. AI ≠ верхний слой данных", {"size": 13, "bold": True, "color": MID}),
        ("Софт работает. Аппаратура отказала. AI зависит от верхнего слоя.", {"size": 11, "color": DEEP}),
    ], line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Урок: регуляторная инфраструктура требует диверсифицированной сенсорной сети, не одного спутника. AI зависит от аппаратуры на верхнем слое.",
                 size=12)
    add_notes(slide, "См. slides/s23-methanesat-loss.md speaker notes.")


def s24_post_methanesat_players(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "После MethaneSAT: Carbon Mapper + GHGSat + Bridger + SeekOps",
             size=24, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Диверсифицированная сенсорная сеть — то, чем должен был быть MethaneSAT. Несколько провайдеров после потери.",
             size=13, italic=True, color=LIGHT)
    players = [
        ("Carbon Mapper Tanager-1", "Запуск 16 авг 2024.\nГруппировка Planet Labs.\nТехнология NASA JPL.\nДетекция на уровне площадки.", TEAL),
        ("Группировка GHGSat", "13 спутников к середине 2025.\nРазрешение 25 м.\nФокус на уровне площадки.\nКоммерческие: ExxonMobil, ConocoPhillips, EOG.", LIGHT),
        ("Авиасъёмка Bridger Photonics", "Воздушный LiDAR (картирование газов).\n4× точнее наземной OGI.\nБлижняя дистанция, медленное сканирование.\nДля программ LDAR.", MID),
        ("SeekOps + Project Canary", "SeekOps — метановые сенсоры на БПЛА.\nProject Canary — непрерывный мониторинг площадок.\nСертификация (RSG, MiQ).", GOLD),
    ]
    p_w = 6.0
    p_h = 2.05
    gap = 0.2
    x0 = 0.5
    y0 = 1.85
    for i, (name, body, accent) in enumerate(players):
        col = i % 2
        row = i // 2
        x = x0 + col * (p_w + gap)
        y = y0 + row * (p_h + gap)
        rounded_box(slide, x, y, p_w, p_h, stroke=accent, stroke_w=2)
        rectangle(slide, x, y, p_w, 0.5, fill=accent)
        text_box(slide, x+0.15, y, p_w-0.3, 0.5, name,
                 size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x+0.2, y+0.6, p_w-0.4, p_h - 0.75, body,
                 size=11, color=DEEP, line_spacing=1.3)
    gold_callout(slide, 0.5, 6.15, 12.33, 0.85,
                 "Урок: один спутник — единая точка отказа. Диверсифицированная сенсорная сеть (спутник + авиа + наземная) — устойчива. Регулятор EU требует именно этого.",
                 size=12)
    add_notes(slide, "См. slides/s24-post-methanesat-players.md speaker notes.")


def s25_4x_discrepancy(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "4× разрыв отрасль vs регулятор — методологический разрыв",
             size=24, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "MethaneSAT измерил ~15 млн т/год по нефтегазу США. EPA Inventory ~4 млн т. Stanford 2024 авиа ~7 млн т = в 2 раза.",
             size=13, italic=True, color=LIGHT)
    rounded_box(slide, 0.5, 1.85, 6.0, 4.5)
    img = ASSETS / "charts" / "s25-4x-discrepancy.png"
    add_image_aspect(slide, img, 0.7, 2.0, 5.6, 4.2)
    rounded_box(slide, 6.7, 1.85, 6.13, 4.5)
    multiline_box(slide, 6.9, 1.95, 5.83, 4.3, [
        ("Что это значит:", {"size": 14, "bold": True, "color": MID}),
        ("· EPA inventory: подход «снизу вверх», коэффициенты 1990-х.", {"size": 11, "color": DEEP}),
        ("· Stanford авиа: измерения «сверху вниз», 2024.", {"size": 11, "color": DEEP}),
        ("· MethaneSAT: космос, всё что меньше 500 кг/ч пропускается.", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("9-спутниковый тест 2024 (Atmospheric Measurement Techniques):", {"size": 13, "bold": True, "color": MID}),
        ("· 58% точек эмиссии идентифицировано.", {"size": 11, "color": DEEP}),
        ("· 41% ложных пропусков.", {"size": 11, "color": DEEP, "bold": True}),
        ("· AI не «ошибается» — сенсоры имеют структурные ограничения.", {"size": 11, "color": DEEP}),
        ("", {"size": 6}),
        ("Это не провал AI:", {"size": 13, "bold": True, "color": GOLD}),
        ("Это структурный методологический разрыв. Модернизация EPA в процессе.", {"size": 11, "color": DEEP, "italic": True}),
    ], line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Триангуляция из нескольких источников — единственный валидный подход. AI необходим, но AI ≠ эталон — он измеряет, факелы факелуют.",
                 size=12)
    add_notes(slide, "См. slides/s25-4x-discrepancy.md speaker notes.")


def s26_eu_vs_epa(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "EU 2024/1787 vs EPA Subpart W — рынок AI MRV развивается асимметрично",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.35, 12.33, 0.4,
             "Оператор в EU рискует штрафом > прибыли от добычи. Оператор в США не торопится инвестировать.",
             size=13, italic=True, color=LIGHT)
    # EU vs EPA boxes
    rounded_box(slide, 0.5, 1.9, 6.0, 4.4)
    rectangle(slide, 0.5, 1.9, 6.0, 0.7, fill=MID)
    text_box(slide, 0.65, 1.9, 5.7, 0.7, "EU 2024/1787 (август 2024)",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 0.7, 2.75, 5.6, 3.5, [
        ("до 20%", {"size": 36, "bold": True, "color": GOLD}),
        ("оборотного штрафа за несоответствие", {"size": 11, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("Ключевые сроки:", {"size": 13, "bold": True, "color": DEEP}),
        ("· LDAR обязательный — 5 мая 2025", {"size": 12, "color": DEEP}),
        ("· OGMP уровень 4/5 отчёты — 5 авг 2025", {"size": 12, "color": DEEP}),
        ("· Соответствие для импорта — с 2027", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Драйвер:", {"size": 12, "bold": True, "color": MID}),
        ("Shell, TotalEnergies — на переднем крае выполнения.", {"size": 12, "color": DEEP}),
    ], line_spacing=1.25)
    rounded_box(slide, 6.7, 1.9, 6.13, 4.4)
    rectangle(slide, 6.7, 1.9, 6.13, 0.7, fill=LIGHT)
    text_box(slide, 6.85, 1.9, 5.83, 0.7, "US EPA Subpart W",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    multiline_box(slide, 6.9, 2.75, 5.83, 3.5, [
        ("6 мая 2024 утверждён → отсрочка до 2034", {"size": 14, "bold": True, "color": DEEP}),
        ("", {"size": 6}),
        ("Что произошло:", {"size": 13, "bold": True, "color": LIGHT}),
        ("· Финальное правило утверждено 6 мая 2024", {"size": 12, "color": DEEP}),
        ("· Администрация Трампа — пересмотр сент. 2025", {"size": 12, "color": DEEP}),
        ("· Предлагается отсрочка до 2034 для внедрения", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Эффект на рынок:", {"size": 13, "bold": True, "color": LIGHT}),
        ("· Рынок AI MRV США ~$200 млн (2024)", {"size": 12, "color": DEEP}),
        ("· vs EU ~$500 млн+ к 2026 (оценка)", {"size": 12, "color": DEEP, "bold": True}),
    ], line_spacing=1.25)
    gold_callout(slide, 0.5, 6.5, 12.33, 0.55,
                 "Регуляторика — главный драйвер Q2. EU жёстко → AI необходим. США в режиме ожидания → рынок AI развивается медленнее.",
                 size=12)
    add_notes(slide, "См. slides/s26-eu-vs-epa-regulation.md speaker notes.")


def s27_q2_alternatives(p):
    slide = blank(p)
    set_slide_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.85,
             "Альтернатива Q2: наземная OGI + переносные анализаторы",
             size=26, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.3, 12.33, 0.4,
             "Когда AI не нужен — OGMP уровень 5 (прямое измерение) + коммерческий учёт нефти.",
             size=13, italic=True, color=LIGHT)
    tools = [
        ("FLIR GFx320 / Opgal EyeCGas", "Ручные OGI-камеры.\nИК-визуализация углеводородов.\n2-3 обхода объекта в год.", LIGHT),
        ("Picarro G2210-i / LI-COR LI-7810", "Переносная спектроскопия затухания в полости.\nПрямое измерение, уровень частей на млрд (ppb).\nДля верификации OGMP уровень 5.", MID),
        ("Rebellion Photonics", "Стационарная OGI + аналитика.\nНепрерывный мониторинг площадки.\nЗаменяет AI-скрининг на крупных узлах.", TEAL),
        ("EPA Method 21 / EU LDAR", "Детектор + газоанализатор для всех клапанов, фланцев.\n4×/год обходы.\nОбязательно по регуляции.", GOLD),
    ]
    t_w = 6.0
    t_h = 2.0
    gap = 0.2
    x0 = 0.5
    y0 = 1.85
    for i, (name, body, accent) in enumerate(tools):
        col = i % 2
        row = i // 2
        x = x0 + col * (t_w + gap)
        y = y0 + row * (t_h + gap)
        rounded_box(slide, x, y, t_w, t_h, stroke=accent, stroke_w=2)
        rectangle(slide, x, y, 0.15, t_h, fill=accent)
        text_box(slide, x + 0.3, y + 0.1, t_w - 0.4, 0.5, name,
                 size=14, bold=True, color=DEEP)
        text_box(slide, x + 0.3, y + 0.65, t_w - 0.4, t_h - 0.75, body,
                 size=11, color=DEEP, line_spacing=1.35)
    gold_callout(slide, 0.5, 6.15, 12.33, 0.85,
                 "2 критерия когда AI не нужен в Q2: (1) соответствие OGMP уровень 5 — регулятор требует точность, не оценку; (2) коммерческий учёт нефти — расходомер обязателен, не AI.",
                 size=12)
    add_notes(slide, "См. slides/s27-q2-alternatives.md speaker notes.")


# Sectionовая загрузка продолжается в build_lec16_p2.py
if __name__ == "__main__":
    print("This is part 1 — see build_all.py for full assembly.")
