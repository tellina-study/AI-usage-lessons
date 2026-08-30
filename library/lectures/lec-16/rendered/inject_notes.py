#!/usr/bin/env python3
"""
inject_notes.py — Phase 8b P0 fix per Лекция 13 lessons.

Извлекает полный текст speaker notes из slides/sNN-*.md
(section after `## Speaker notes`) и инжектирует в PPTX через
notes_slide.notes_text_frame.

Запускать ПОСЛЕ build_lec16_p2.py — каждый rebuild стирает notes.

Pipeline:
  python3 build_lec16_p2.py
  python3 inject_notes.py

Issue #144.
"""
import re
import sys
from pathlib import Path

from pptx import Presentation

WT = Path('/tmp/lec-16-wt/library/lectures/lec-16')
PPTX = WT / 'rendered' / 'lec-16.pptx'
SLIDES_DIR = WT / 'slides'


def extract_speaker_notes(md_path: Path) -> str | None:
    """Извлечь текст после `## Speaker notes` до следующего `## ` или EOF."""
    text = md_path.read_text(encoding='utf-8')
    m = re.search(r'^##\s+Speaker notes\s*$', text, re.MULTILINE)
    if not m:
        m = re.search(r'^###\s+Speaker notes\s*$', text, re.MULTILINE)
    if not m:
        return None
    after = text[m.end():].strip()
    # отрезать следующую секцию того же уровня (`## ` без `###`)
    next_h = re.search(r'^##\s+(?!#)', after, re.MULTILINE)
    if next_h:
        after = after[:next_h.start()].strip()
    return after


def main() -> int:
    if not PPTX.exists():
        print(f'ERROR: PPTX not found: {PPTX}', file=sys.stderr)
        return 1

    p = Presentation(PPTX)
    slide_files = sorted(SLIDES_DIR.glob('s*.md'))
    print(f'Found {len(slide_files)} markdown slides; deck has {len(p.slides)}')

    if len(slide_files) != len(p.slides):
        print('WARNING: count mismatch — proceeding with zip-min', file=sys.stderr)

    injected = 0
    for idx, (slide, md) in enumerate(zip(p.slides, slide_files), 1):
        notes_text = extract_speaker_notes(md)
        if notes_text is None:
            print(f's{idx:02d}: NO speaker notes section in {md.name} — skipping')
            continue
        nf = slide.notes_slide.notes_text_frame
        nf.clear()
        nf.paragraphs[0].text = notes_text
        print(f's{idx:02d}: injected {len(notes_text)} chars from {md.name}')
        injected += 1

    p.save(PPTX)
    print(f'\nPPTX saved: {injected}/{len(p.slides)} slides got full notes injected')
    return 0


if __name__ == '__main__':
    sys.exit(main())
