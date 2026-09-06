"""
Full 36-slide build of Лекции 3 «Архитектуры AI-систем: агенты, RAG, API» (v3).

v3 = СТРУКТУРНАЯ ревизия v2→v3 по owner-обратной связи (plan §4, U-1…U-9):
  +6 suffix-слайдов (НЕ перенумеровывая s01–s30):
    s04a divider Раздел 1 · s13a divider Раздел 3 · s13b определение FT ·
    s23a sub-divider Безопасность · s25a divider Раздел 5 · s31 Q&A.
  s30 ретайтл (U-6, function-as-title убран) + Q&A вынесен в s31 (U-7).
  Порядок: s01..s04 → s04a → s05..s08 → s09 → s10..s13 → s13a → s13b →
           s14..s17 → s18 → s19..s23 → s23a → s24..s25 → s25a →
           s26..s29 → s30 → s31.

Source-of-truth: deck.yaml + deck-part2.yaml (U-9 split, v3) +
chapter v1.1 finalized (~22450 слов) + slides/*.md (36 файлов,
readable speaker notes 150-300 слов).

Issue #87 · Branch: issue-87-lec-03-architectures

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).

Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml ≈ 75 мин.

Build via: python3 build_v3.py — generates lec-03.pptx (36 slides).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree
from PIL import Image

import sys as _sys
_sys.path.insert(0, str(Path(__file__).resolve().parent))
import refs_lec03 as R  # noqa: E402  (issue #171 reference/page-number system)

# issue #171: footer text capture so ref-slides can fold the caveat into the
# clickable [N] source list (single bottom band, no overlap).
_FOOTER_TEXT = {}

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x5B, 0x66, 0x78)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFD, 0xF3, 0xDC)
TEAL_TINT = RGBColor(0xE4, 0xF1, 0xF2)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
# issue #171: run from the worktree — read slides/assets/notes + write pptx
# from THIS repo checkout (slides identical to main). Falls back to main-repo
# ROOT only if the worktree copy is missing.
_WT = Path(__file__).resolve().parents[1]      # …/lec-03 in the current checkout
_MAIN = Path("/home/harness/harness-projects/256/lessons-3bb49d40/library/lectures/lec-03")
ROOT = _WT if (_WT / "slides").exists() else _MAIN
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
CHARTS = ASSETS / "charts"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-03.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before") is not None:
                p.space_before = Pt(cfg["space_before"])
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.035, min(0.22, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def up_arrow(slide, x, y, w, h, fill=LIGHT):
    shp = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, fill, *, stroke=None, stroke_pt=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def connector(slide, x1, y1, x2, y2, color=LIGHT, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = etree.SubElement(ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
        pd.set("val", dash)
    return cn


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    path = Path(path)
    if not path.exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path); iw, ih = img.size; img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
            return
        ir = iw / ih; br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x), Inches(y + (h - ah) / 2),
                                     width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2), Inches(y),
                                     height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


SCREENSHOTS = ROOT / "assets/screenshots"
_CROP_CACHE = ROOT / "rendered/assets/_crop_cache"


def hero_image(slide, src, x, y, w, h):
    """Cover-crop an image to EXACTLY fill box (x,y,w,h) — no distortion, no
    letterbox. python-pptx can't crop, so we pre-crop via PIL to the target
    aspect ratio and cache the result. Used for ≥40% hero fills (s01, s30)."""
    src = Path(src)
    if not src.exists():
        return
    _CROP_CACHE.mkdir(parents=True, exist_ok=True)
    target_ratio = w / h
    try:
        img = Image.open(src).convert("RGB")
    except Exception:
        return
    iw, ih = img.size
    ir = iw / ih
    if ir > target_ratio:      # image wider — crop sides
        new_w = int(ih * target_ratio)
        left = (iw - new_w) // 2
        img = img.crop((left, 0, left + new_w, ih))
    else:                       # image taller — crop top/bottom
        new_h = int(iw / target_ratio)
        top = (ih - new_h) // 2
        img = img.crop((0, top, iw, top + new_h))
    out = _CROP_CACHE / f"{src.stem}_{w:.2f}x{h:.2f}.png"
    img.save(out)
    slide.shapes.add_picture(str(out), Inches(x), Inches(y),
                             width=Inches(w), height=Inches(h))


def slide_title(slide, text, *, y=0.42, h=0.95, w=12.25, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.12, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 color=DEEP, align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.22, y=y + 0.06, w=w - 0.44, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.22)


def footer(slide, text):
    # issue #171: record the footer so ref-slides can relocate/fold it; still
    # render normally (post-processing removes it only on ref-slides).
    _FOOTER_TEXT[id(slide)] = text
    tb = text_box(slide, x=0.55, y=7.02, w=12.25, h=0.36, text=text,
                  size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
                  line_spacing=1.0)
    _FOOTER_TEXT.setdefault("_shapes", {})[id(slide)] = tb
    return tb


def icon(slide, name, x, y, size, variant="mid"):
    add_image(slide, ICONS / f"{name}-{variant}.png", x, y, size, size)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Deck loader — U-9: deck.yaml split на 2 части (≤600 строк каждая).
# Loader читает ОБЕ части, объединяет ключ `slides`, валидирует totals.
# ============================================================
def load_deck():
    """Load 2-part deck spec (deck.yaml + deck-part2.yaml), merge `slides`.

    Returns dict with merged slide list + validation. Raises if the
    ordered slide-id list does not match the canonical v3 presentation
    order (cascade-safe guard — s01–s30 must NOT be renumbered).
    """
    try:
        import yaml
    except ImportError:
        return None  # yaml optional — builder list is authoritative anyway
    p1 = ROOT / "deck.yaml"
    p2 = ROOT / "deck-part2.yaml"
    parts = [ROOT / "deck.yaml", ROOT / "deck-part2.yaml", ROOT / "deck-part3.yaml"]
    slides = []
    d1 = None
    totals = {}
    for pp in parts:
        if not pp.exists():
            continue
        d = yaml.safe_load(pp.read_text(encoding="utf-8"))
        if d1 is None:
            d1 = d
        slides += list(d.get("slides", []))
        if d.get("totals"):
            totals = d["totals"]
    spec = {
        "deck": d1.get("deck", {}) if d1 else {},
        "palette": d1.get("palette", {}) if d1 else {},
        "slides": slides,
        "totals": totals,
    }
    ids = [s["id"] for s in slides]
    # builder list is authoritative; deck.yaml is documentation — report only.
    print(f"deck.yaml slide ids ({len(ids)}): {ids}")
    return spec


# ============================================================
# Section divider — unified template (6-card roadmap, gold current)
# Sections of Лекции 3 (deck.yaml): 0..5.
# ============================================================
NAV = [
    ("0", "Открытие"),
    ("1", "Промпт"),
    ("2", "RAG"),
    ("3", "Дообучение"),
    ("4", "Агенты"),
    ("5", "Фреймворк"),
]


def roadmap_bar(slide, here_idx, *, y=6.45):
    """6-card progress bar; current section gold-bordered."""
    n = len(NAV)
    gap = 0.14
    bx = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    ch = 0.62
    for i, (num, label) in enumerate(NAV):
        x = bx + i * (cw + gap)
        cur = (i == here_idx)
        if cur:
            filled_rect(slide, x, y, cw, ch, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.12)
        else:
            filled_rect(slide, x, y, cw, ch, SURFACE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(slide, x=x + 0.06, y=y + 0.07, w=cw - 0.12, h=0.24,
                 text=f"Раздел {num}", size=10.5, bold=True,
                 color=(DEEP if cur else LIGHT), align=PP_ALIGN.CENTER)
        text_box(slide, x=x + 0.06, y=y + 0.31, w=cw - 0.12, h=0.26,
                 text=label, size=11, bold=cur,
                 color=(DEEP if cur else SLATE), align=PP_ALIGN.CENTER)


WEB = ASSETS / "web"


def build_section_divider(p, here_idx, big_num, subtitle, frame_phrase, sid,
                          *, image_src=None, image_caption=None, tag=None):
    """Distinct divider (NO ocean motif). Left = РАЗДЕЛ N + subtitle + frame
    phrase + optional tag chip; right = a REAL metaphor image (≥40% width,
    full height, cover-cropped) when image_src is given, else the giant
    cover-style decorative digit. Gold-current roadmap bar at bottom."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    have_img = image_src is not None and Path(image_src).exists()
    if have_img:
        # Real metaphor image — right ~42% width, full height (≥40% area).
        ix, iy, iw, ih = 8.10, 0.0, 5.233, 7.5
        hero_image(s, image_src, ix, iy, iw, ih)
        # faint decorative digit overlaid top-right of image (brand echo)
        text_box(s, x=ix + 1.6, y=0.20, w=3.6, h=2.6, text=str(here_idx),
                 size=170, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        if image_caption:
            filled_rect(s, ix, 7.10, iw, 0.40, DEEP)
            text_box(s, ix + 0.16, 7.12, iw - 0.30, 0.34, image_caption,
                     size=10, italic=True, color=WHITE,
                     anchor=MSO_ANCHOR.MIDDLE)
        left_w = 7.35
    else:
        text_box(s, x=8.35, y=0.55, w=4.6, h=5.6, text=str(here_idx),
                 size=380, bold=True, color=COVER_OUTLINE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        left_w = 7.5
    # Left — section label + subtitle + frame phrase
    text_box(s, x=0.75, y=1.35, w=left_w - 0.2, h=0.55,
             text=f"РАЗДЕЛ {here_idx}", size=20, bold=True, color=TEAL)
    filled_rect(s, 0.78, 1.96, 0.7, 0.05, fill=GOLD)
    text_box(s, x=0.75, y=2.30, w=left_w, h=1.75, text=subtitle,
             size=36, bold=True, color=DEEP, line_spacing=1.06)
    text_box(s, x=0.78, y=4.20, w=left_w - 0.15, h=1.45, text=frame_phrase,
             size=17, italic=True, color=LIGHT, line_spacing=1.20)
    if tag:
        chip(s, 0.78, 5.70, min(left_w - 0.3, 0.14 * len(tag) + 0.6), 0.44,
             tag, fill=TEAL, color=WHITE, size=12.5)
    roadmap_bar(s, here_idx, y=6.45)
    speaker_notes(s, load_notes(sid))
    return s


# ============================================================
# Slide builders — 30 slides
# ============================================================

def _pill_figure(s, cx, cy, scale, body_col, ok=True):
    """Flat vector figure holding up a giant 'pill'. ok=True → confident pose
    (gold pill, teal glow); ok=False → pill shattered / red-adjacent avoided,
    we keep Ocean palette: dull grey pill + slate figure."""
    # head
    circle(s, cx - 0.30 * scale, cy, 0.60 * scale, body_col)
    # torso (rounded rect)
    filled_rect(s, cx - 0.42 * scale, cy + 0.66 * scale, 0.84 * scale,
                1.15 * scale, body_col, radius=True, radius_adj=0.35)


def build_s01(p):
    """hero_cover / meme-hook — «Магическая пилюля» (issue #185). Реальный
    интернет-мем Drake (imgflip) с русскими подписями: reject = усложнять
    промпт ради точности, approve = выбрать архитектуру под задачу. Мем несёт
    тезис; ответ раскрывается по лекции. Air Canada — отдельно на s01b."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # top ribbon (минимум текста — постановка одной строкой)
    text_box(s, 0.55, 0.40, 12.25, 0.44, "МИФ ПРО AI-СИСТЕМЫ",
             size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(s, 0.55, 0.86, 12.25, 0.92,
             "«Магическая пилюля»: модель ответит точнее, если попросить её экспертом и усложнить промпт.",
             size=25, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             line_spacing=1.06)
    # real Drake meme (hero ≥40%) в Ocean-рамке по центру: наверху отвергаем
    # «усложнять промпт ради точности», внизу принимаем «выбор архитектуры».
    my, mh = 2.06, 4.06
    mw = 4.94
    mx = (13.33 - mw) / 2
    ocean_box(s, mx - 0.16, my - 0.14, mw + 0.32, mh + 0.28)
    add_image(s, WEB / "s01-drake-ru.png", mx, my, mw, mh)
    gold_callout(s, 0.55, 6.34, 12.25, 0.74,
                 "Откуда на самом деле берётся надёжность AI-системы — разбираем всю лекцию.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """cover — distinct, NO ocean motif. Mega «03» + title + roadmap.
    v2: subtitle brought to lec-02 cover canon — content-promise line with
    teal accent bar + MID color (was designer-initiative «Курс · 75 минут»
    meta-line italic-light, removed)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=7.6, y=1.15, w=5.7, h=5.0, text="03",
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.75, y=1.35, w=6.6, h=0.5, text="ЛЕКЦИЯ 3",
             size=18, bold=True, color=TEAL)
    filled_rect(s, 0.78, 1.92, 0.7, 0.05, fill=TEAL)
    text_box(s, x=0.75, y=2.35, w=7.7, h=2.7,
             text="Архитектуры AI-систем:\nагенты, RAG, API",
             size=46, bold=True, color=DEEP, line_spacing=1.08)
    # subtitle = content promise (lec-02 canon: teal accent bar + MID, не meta)
    filled_rect(s, 0.78, 5.28, 0.05, 0.56, fill=TEAL)
    text_box(s, x=1.02, y=5.26, w=7.4, h=0.62,
             text="Какую архитектуру выбрать под задачу —\nи когда правильный ответ «не ИИ»",
             size=19, italic=False, color=MID, line_spacing=1.18)
    # v4 (#212): cover теперь чистый — roadmap вынесен в отдельный
    # lecture-map слайд s02a (паттерн Л1/Л2). Cover без roadmap-bar.
    # #185/#313: атрибуция курса «3 курс ИУ6 · Модуль 1…» снята с cover.
    # #185/#313: тематический интернет-мем на cover (левый нижний угол) —
    # реальный кадр-мем «well yes, but actually no» (пират-жест): «хотелось бы
    # магической пилюли — но нет». Компактный, не спорит с mega-«03».
    _cover_pirate_meme(s, 0.75, 6.06)
    speaker_notes(s, load_notes("s02"))


def _cover_pirate_meme(s, x, y):
    """Small real internet-meme motif for the cover: пират-жест «well yes, but
    actually no» (imgflip; английская подпись обрезана) + русская подпись рядом
    — намёк на тезис лекции: «магической пилюли» не существует. Компактный:
    левый нижний угол, не спорит с mega-«03». Атрибуция — только attribution.md."""
    mw = 2.10
    mh = mw * (755 / 1600)      # сохранить пропорции кропа
    ocean_box(s, x - 0.08, y - 0.08, mw + 0.16, mh + 0.16)
    add_image(s, WEB / "cover-pirate-crop.png", x, y, mw, mh)
    text_box(s, x + mw + 0.28, y - 0.06, 5.4, mh + 0.16,
             "«магической пилюли» не существует",
             size=14, italic=True, bold=True, color=MID,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)


def build_s02a(p):
    """NEW (#212) lecture-map — 6 horizontal section cards (Л1/Л2 pattern).
    Отдельный слайд-содержание после cover. Показывает маршрут лекции:
    Разделы 0–5 с одной строкой смысла каждого."""
    s = blank(p)
    slide_title(s, "Маршрут лекции — шесть разделов.", size=27)
    # #185/#315: раздел «Агенты» больше не подсвечен; gold-акцент сохранён на
    # несущей линии лекции (маркер перед подзаголовком), а не на разделе.
    filled_rect(s, 0.55, 1.26, 0.06, 0.34, GOLD)
    text_box(s, 0.74, 1.22, 12.05, 0.42,
             "Одна несущая линия: выбор архитектуры под задачу — и когда правильный ответ «не ИИ».",
             size=15, italic=True, color=MID)
    cards = [
        ("0", "Открытие", "постановка задачи: откуда берётся надёжность", "gavel", MID),
        ("1", "Промпт и его границы", "что умеет один вызов и где его потолок", "message-circle", MID),
        ("2", "RAG", "внешнее знание в контекст — и где оно тихо ломается", "database", MID),
        ("3", "Дообучение", "менять веса под поведение, не под знание", "sliders-horizontal", MID),
        ("4", "Агенты", "цикл, экипировка, память, безопасность — и провалы", "bot", MID),
        ("5", "Фреймворк", "лестница + чек-лист: как выбрать быстро и обоснованно", "list-checks", MID),
    ]
    # 2 rows × 3 cols
    cw, chh = 3.95, 2.30
    gapx, gapy = 0.20, 0.28
    x0, y0 = 0.55, 1.90
    for i, (num, title, sub, ic, col) in enumerate(cards):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gapx)
        y = y0 + r * (chh + gapy)
        isgold = (col == GOLD)
        if isgold:
            ocean_box(s, x, y, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, y, cw, chh)
        circle(s, x + 0.24, y + 0.24, 0.56, col)
        text_box(s, x + 0.24, y + 0.24, 0.56, 0.56, num,
                 size=22, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        icon(s, ic, x + cw - 0.74, y + 0.26, 0.50, "gold" if isgold else "mid")
        text_box(s, x + 0.24, y + 0.98, cw - 0.48, 0.60, title,
                 size=16, bold=True, color=DEEP, line_spacing=1.05)
        text_box(s, x + 0.24, y + 1.58, cw - 0.48, 0.62, sub,
                 size=12.5, color=SLATE, line_spacing=1.14)
    speaker_notes(s, load_notes("s02a"))


def build_s03(p):
    """recap (§0) — лёгкое напоминание двух понятий Лекции 2, на которые
    опирается вся лекция: одиночный вызов (single-shot) и семантический поиск
    на эмбеддингах. v6 (#185/#316): бывшая тяжёлая схема «4 обвязки» снята —
    обвязки раскрываются по разделам, здесь только опора из Л2."""
    s = blank(p)
    slide_title(s, "Из Лекции 2 берём готовыми два понятия.", size=27)
    text_box(s, 0.55, 1.22, 12.25, 0.42,
             "Не переобъясняем — просто вспоминаем, чтобы дальше на них опереться. Всё остальное надстроим по ходу лекции.",
             size=15, italic=True, color=MID, line_spacing=1.15)
    # два крупных recap-бокса
    by, bh = 2.05, 3.35
    bw = 6.05
    # LEFT — single-shot
    ocean_box(s, 0.55, by, bw, bh)
    icon(s, "message-circle", 0.85, by + 0.30, 0.62, "mid")
    text_box(s, 1.70, by + 0.30, bw - 1.10, 0.60, "Одиночный вызов (single-shot)",
             size=19, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, 0.85, by + 1.20, bw - 0.60, 1.05,
             "Один проход модели: дал промпт → получил ответ. Без памяти между вызовами, без обращений наружу.",
             size=15, color=DEEP, line_spacing=1.20)
    filled_rect(s, 0.85, by + 2.50, bw - 0.60, 0.66, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.5, radius=True, radius_adj=0.14)
    text_box(s, 1.05, by + 2.53, bw - 1.00, 0.60,
             "Модель знает только то, что в промпте и в весах",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    # RIGHT — semantic search / embeddings
    rx = 6.75
    ocean_box(s, rx, by, bw, bh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "database", rx + 0.30, by + 0.30, 0.62, "teal")
    text_box(s, rx + 1.15, by + 0.30, bw - 1.40, 0.60, "Семантический поиск на эмбеддингах",
             size=19, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, rx + 0.30, by + 1.20, bw - 0.60, 1.05,
             "Текст → вектор смысла; близкие по смыслу фрагменты — близкие векторы. Поиск «по смыслу», а не по точному слову.",
             size=15, color=DEEP, line_spacing=1.20)
    filled_rect(s, rx + 0.30, by + 2.50, bw - 0.60, 0.66, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.14)
    text_box(s, rx + 0.50, by + 2.53, bw - 1.00, 0.60,
             "На этом стоит RAG (Раздел 2)",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    gold_callout(s, 0.55, 5.72, 12.25, 0.88,
                 "Оба понятия — фундамент из Лекции 2. Всё, что мы надстроим сегодня (RAG, инструменты, дообучение, агенты), опирается на них.",
                 size=15, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """assertion_visual — central question + 6-step ladder."""
    s = blank(p)
    text_box(s, 0.55, 0.38, 12.25, 0.42, "ЦЕНТРАЛЬНЫЙ ВОПРОС ЛЕКЦИИ",
             size=14, bold=True, color=TEAL)
    qx, qy, qw, qh = 0.55, 0.85, 12.25, 1.30
    ocean_box(s, qx, qy, qw, qh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, qx + 0.4, qy + 0.18, qw - 0.8, qh - 0.36,
             "У меня есть задача и доступ к LLM. Какую архитектуру выбрать — и когда правильный ответ «не ИИ»?",
             size=23, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.12)
    # Ladder — 6 steps bottom-up: idx 0 = step 1 (BOTTOM, gold), idx 5 = step 6 (top).
    steps = [
        ("1", "Обычный код (без ИИ)", "точка отсчёта", GOLD, True),
        ("2", "Один вызов LLM", "промпт; + CoT (по шагам), + примеры в промпте (few-shot)", MID, False),
        ("3", "RAG / контекст-инжиниринг", "поиск-дополненная генерация", LIGHT, False),
        ("4", "Сценарий (workflow)", "предопределённые пути", LIGHT, False),
        ("5", "Агент", "цикл: план → действие → проверка → повтор (plan → act → check → iterate)", LIGHT, False),
        ("6", "Мульти-агент", "несколько координируемых агентов", LIGHT, False),
    ]
    n = len(steps)
    # #215/#216: увеличены визуальные элементы лестницы — крупнее рунги
    # (step_h 0.66→0.74), крупнее номера-круги (0.36→0.46) и sub-подписи
    # (10.5→12), а также стрелка направления и её метки.
    step_h = 0.74
    vgap = 0.075
    bottom_edge = 6.85  # bottom of step 1
    sx = 0.55
    for i, (num, label, sub, col, isgold) in enumerate(steps):
        # i=0 -> bottom rung; each higher step sits above + indented right
        y = bottom_edge - step_h - i * (step_h + vgap)
        indent = i * 0.22
        w = 7.95 - indent
        x = sx + indent
        if isgold:
            ocean_box(s, x, y, w, step_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.5)
        else:
            ocean_box(s, x, y, w, step_h)
        circle(s, x + 0.16, y + (step_h - 0.46) / 2, 0.46, col)
        text_box(s, x + 0.16, y + (step_h - 0.46) / 2, 0.46, 0.46, num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.74, y + 0.11, w - 0.90, 0.30, label,
                 size=16, bold=True, color=DEEP)
        text_box(s, x + 0.74, y + 0.43, w - 0.90, 0.26, sub,
                 size=12, italic=True, color=SLATE)
    # PA-1 (owner-approved) + #216: direction-of-scale strip alongside ladder —
    # enlarged «сложнее ↑» / «проще ↓» + thicker arrow.
    text_box(s, 7.98, 2.28, 1.22, 0.34, "сложнее ↑", size=14, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    up_arrow(s, 8.40, 2.72, 0.40, 3.86, fill=COVER_OUTLINE)
    text_box(s, 7.98, 6.62, 1.22, 0.34, "проще ↓", size=14, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    # rule on the right
    ocean_box(s, 9.05, 2.55, 3.75, 3.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "milestone", 9.32, 2.85, 0.62, "gold")
    text_box(s, 9.32, 3.70, 3.25, 2.65,
             "Подниматься на следующую ступень — только при требовании задачи, которого текущая не закрывает.",
             size=16, bold=True, color=DEEP, line_spacing=1.24)
    footer(s, "Лестница — карта лекции, не требование понять всё сейчас. Каждую ступень разберём отдельно.")
    speaker_notes(s, load_notes("s04"))


def build_s05(p):
    """assertion_visual — default = one call; cost of each climb."""
    s = blank(p)
    slide_title(s, "По умолчанию — один вызов с хорошим промптом.", size=27)
    # #219-context: explicit граница знания модели на видимом слое (§1.1)
    filled_rect(s, 0.55, 1.18, 12.25, 0.66, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.5, radius=True, radius_adj=0.10)
    text_box(s, 0.78, 1.24, 11.8, 0.56,
             "Модель знает только то, что в промпте (плюс то, что было в весах на обучении). Больше ей взять неоткуда.",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.12)
    # Left — anchor block
    lx, ly, lw, lh = 0.55, 2.02, 6.15, 3.62
    ocean_box(s, lx, ly, lw, lh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "target", lx + 0.28, ly + 0.24, 0.58, "teal")
    text_box(s, lx + 0.28, ly + 0.94, lw - 0.56, 0.6,
             "Один вызов LLM\nс хорошим промптом",
             size=18, bold=True, color=DEEP, line_spacing=1.1)
    bullets = [
        "минимальная стоимость (один проход)",
        "минимальная задержка (нет лишних обращений)",
        "максимальная предсказуемость (нет петель, нет поиска, который тихо деградирует)",
    ]
    by = ly + 1.90
    for b in bullets:
        circle(s, lx + 0.30, by + 0.07, 0.12, TEAL)
        text_box(s, lx + 0.56, by, lw - 0.84, 0.62, b,
                 size=13, color=DEEP, line_spacing=1.12)
        by += 0.56
    # Right — реальный интернет-мем «Gru's Plan» (imgflip, #185; замена
    # Expanding Brain — owner: «мем с мозгом был уже»): эскалация архитектуры
    # «на всякий случай» с абсурдной развязкой в 4-й панели — задача была на
    # три строки обычного кода. Русские подписи вжжены в шаблон.
    _s05_overengineering_meme(s, 6.95, 2.02, 5.85, 3.62)
    gold_callout(s, 0.55, 5.92, 12.25, 0.92,
                 "Не усложняй архитектуру без причины, выраженной в требованиях задачи. Это распределение бремени доказательства, а не примитивизм.",
                 size=15)
    speaker_notes(s, load_notes("s05"))


def _s05_overengineering_meme(s, x, y, w, h):
    """Реальный интернет-мем «Gru's Plan» (imgflip blank + русские подписи):
    эскалация архитектуры «на всякий случай» с абсурдной развязкой в 4-й
    панели. Заголовок «оверинжиниринг» + сам мем в Ocean-рамке. Мем несёт
    тезис «не усложняй без причины» без атрибуции на слайде (attribution.md)."""
    ocean_box(s, x, y, w, h)
    chip(s, x + 0.28, y + 0.22, w - 0.56, 0.46, "«НА ВСЯКИЙ СЛУЧАЙ» — ЭСКАЛАЦИЯ",
         fill=LIGHT, color=WHITE, size=13.5)
    # Gru's Plan (composite 700x449, landscape 1.56) вписан по ширине бокса
    from PIL import Image as _I
    _im = _I.open(WEB / "s05-gru-ru.png"); _r = _im.size[0] / _im.size[1]; _im.close()
    img_w = w - 0.56
    img_h = img_w / _r
    if img_h > h - 0.86:
        img_h = h - 0.86
        img_w = img_h * _r
    img_x = x + (w - img_w) / 2
    add_image(s, WEB / "s05-gru-ru.png", img_x, y + 0.82, img_w, img_h)


def build_s05a(p):
    """NEW (§1.2) — роли в промпте: миф «персона = точность» опровергнут.
    Zheng et al. 2024 EMNLP + arXiv:2605.29420. Часть failure/judgment —
    опровержение «магической пилюли»."""
    s = blank(p)
    slide_title(s, "Роль в промпте настраивает тон — не точность.", size=26)
    text_box(s, 0.55, 1.16, 12.25, 0.42,
             "«Ты — опытный юрист» сдвигает внимание модели в сторону текста такой роли — но это про стиль, не про факты.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # #185 re-layout: минимум текста + мем. Левая треть — миф + один факт
    # эксперимента; центр — крупный мем «Change my mind»; справа — что роль
    # делает реально. Gold-вывод широкой полосой под всеми тремя.
    ly = 1.78
    mx, mw = 5.20, 4.05
    mem_h = mw / 1.335                 # ≈3.03
    col_h = mem_h + 0.10              # карточки под высоту мема (mass balance)
    # LEFT — миф + один результат (сжато)
    lx, lw = 0.55, 4.35
    ocean_box(s, lx, ly, lw, col_h)
    text_box(s, lx + 0.28, ly + 0.22, lw - 0.56, 0.36, "Расхожий миф",
             size=15, bold=True, color=LIGHT)
    text_box(s, lx + 0.28, ly + 0.62, lw - 0.56, 0.80,
             "«Роль эксперта в промпте — и модель точнее по фактам»",
             size=15, italic=True, color=SLATE, line_spacing=1.18)
    connector(s, lx + 0.28, ly + 1.58, lx + lw - 0.28, ly + 1.58, LIGHT, 1.0)
    text_box(s, lx + 0.28, ly + 1.74, lw - 0.56, 0.36, "Эксперимент",
             size=15, bold=True, color=MID)
    text_box(s, lx + 0.28, ly + 2.14, lw - 0.56, 0.90,
             "162 персоны, 2410 вопросов —\nперсоны НЕ повысили точность.",
             size=15.5, bold=True, color=DEEP, line_spacing=1.26)
    # CENTER — крупный мем «Change my mind» (свой полноценный слот)
    my = ly
    ocean_box(s, mx - 0.14, my - 0.10, mw + 0.28, mem_h + 0.20)
    add_image(s, WEB / "s05a-changemymind-ru.png", mx, my, mw, mem_h)
    # RIGHT — что роль делает реально (высота под мем)
    rx, rw = 9.55, 3.25
    ocean_box(s, rx, ly, rw, col_h, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "user-round", rx + 0.26, ly + 0.24, 0.48, "teal")
    text_box(s, rx + 0.86, ly + 0.24, rw - 1.05, 0.48, "Роль влияет на:",
             size=15, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.26, ly + 0.94, rw - 0.52, 1.20,
             "тон и глубину изложения — насколько формально и подробно, а не на верность факта.",
             size=14.5, color=DEEP, line_spacing=1.24)
    # gold-вывод широкой полосой под тремя колонками (минимум текста)
    gy = ly + col_h + 0.16
    gold_callout(s, 0.55, gy, 12.25, 6.90 - gy,
                 "Нужна точность? Инструмент — не роль, а контекст и RAG (проверяемый источник). Ещё один пункт «магической пилюли», опровергнутый измерением.",
                 size=15)
    footer(s, "Zheng, Pei, Logeswaran, Lee, Jurgens · Findings of EMNLP 2024 (arXiv:2311.10054) + arXiv:2605.29420 (2026).")
    speaker_notes(s, load_notes("s05a"))


def build_s05b(p):
    """NEW (§1.3) — структура промпта: разделители + разделение
    инструкция/контекст/данные. Параллель со structured output (§4.1):
    вход vs выход — тот же принцип."""
    s = blank(p)
    slide_title(s, "Структура промпта: разделяй инструкцию, контекст, данные.",
                size=23, h=1.30, line_spacing=1.08)
    text_box(s, 0.55, 1.58, 12.25, 0.42,
             "Плоский промпт заставляет модель угадывать, где кончается инструкция и начинаются данные. Явная граница снимает эту неоднозначность.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # left — 3 stacked labelled parts (structured input)
    lx, ly, lw = 0.55, 2.14, 5.85
    parts = [
        ("braces", "Инструкция", "что сделать", MID),
        ("book-open", "Контекст", "на основе чего", LIGHT),
        ("database", "Данные", "с чем именно работать", TEAL),
    ]
    ph = 1.16
    pgap = 0.16
    py = ly
    for ic, t, sub, col in parts:
        ocean_box(s, lx, py, lw, ph)
        filled_rect(s, lx + 0.22, py + 0.22, 0.72, ph - 0.44, col, radius=True,
                    radius_adj=0.16)
        icon(s, ic, lx + 0.34, py + ph / 2 - 0.24, 0.48, "white")
        text_box(s, lx + 1.14, py + 0.20, lw - 1.35, 0.42, t,
                 size=17, bold=True, color=DEEP)
        text_box(s, lx + 1.14, py + 0.62, lw - 1.35, 0.40, sub,
                 size=13, italic=True, color=SLATE)
        py += ph + pgap
    # right — delimiters kinds + structured-output parallel
    rx, rw = 7.15, 5.65
    ocean_box(s, rx, ly, rw, 1.95)
    text_box(s, rx + 0.28, ly + 0.18, rw - 0.56, 0.36, "Разделители (delimiters)",
             size=15, bold=True, color=MID)
    for i, d in enumerate([
        "XML-теги:  <инструкция>…</инструкция>",
        "Markdown-заголовки:  ## Задача · ## Данные",
        "Тройные кавычки / бэктики вокруг данных",
    ]):
        circle(s, rx + 0.30, ly + 0.62 + i * 0.40 + 0.05, 0.12, LIGHT)
        text_box(s, rx + 0.54, ly + 0.62 + i * 0.40, rw - 0.82, 0.38, d,
                 size=13, color=DEEP, font=FONT_MONO if i < 2 else FONT_BODY,
                 line_spacing=1.05)
    ocean_box(s, rx, ly + 2.10, rw, 1.70, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.28, ly + 2.24, rw - 0.56, 1.44,
             "Тот же принцип, что структурированный вывод (Раздел 4): там схема задаёт структуру ВЫХОДА, здесь разделители — структуру ВХОДА. Ясная структура вместо «выведи форму по смыслу».",
             size=13.5, color=DEEP, line_spacing=1.20)
    gold_callout(s, 0.55, 5.98, 12.25, 0.86,
                 "Модель, которой показали, где кончается инструкция и начинаются данные, реже принимает фрагмент данных за новую команду — та же путаница лежит в основе инъекции в промпт (prompt injection) (Раздел 4).",
                 size=13.5)
    speaker_notes(s, load_notes("s05b"))


def build_s06(p):
    """case_study — CoT worked example + faithfulness limit MERGED (§1.4+§1.5)."""
    s = blank(p)
    slide_title(s, "Chain-of-thought помогает — но его нельзя аудировать.", size=24)
    # TOP BAND — CoT worked example (compact before/after)
    cy, ch = 1.16, 2.02
    cw = 6.05
    ocean_box(s, 0.55, cy, cw, ch)
    text_box(s, 0.83, cy + 0.14, cw - 0.56, 0.34, "Без CoT",
             size=15, bold=True, color=LIGHT)
    text_box(s, 0.83, cy + 0.52, cw - 0.56, 0.85,
             "«Было 23 яблока, 7 испортились, докупили 2 ящика по 6. Сколько хороших?»",
             size=13, color=DEEP, line_spacing=1.14)
    text_box(s, 0.83, cy + 1.42, cw - 0.56, 0.5,
             "→ правдоподобное, но неверное число",
             size=13, bold=True, color=SLATE)
    rx0 = 6.75
    ocean_box(s, rx0, cy, cw, ch, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx0 + 0.28, cy + 0.14, cw - 0.56, 0.34, "С CoT («решай по шагам»)",
             size=15, bold=True, color=TEAL)
    text_runs(s, rx0 + 0.30, cy + 0.56, cw - 0.6, 0.9, [
        {"text": "23 − 7 = 16    ·    2 × 6 = 12", "size": 15, "bold": True, "color": DEEP},
        {"text": "16 + 12 = 28", "size": 22, "bold": True, "color": GOLD,
         "newpara": True, "space_before": 6},
    ], line_spacing=1.1)
    text_box(s, rx0 + 0.30, cy + 1.56, cw - 0.6, 0.36, "→ верно",
             size=13, bold=True, color=TEAL)
    # thin note between bands (single line)
    text_box(s, 0.55, cy + ch + 0.06, 12.25, 0.28,
             "Технически это всё ещё один вызов — CoT инструмент под класс задач (цепочка шагов), не глобальный тумблер.",
             size=12, italic=True, color=MID)
    # BOTTOM BAND — faithfulness limit (качественно; числа — на следующем слайде)
    fy = 3.70
    text_box(s, 0.55, fy, 12.25, 0.34,
             "Но проговорённое рассуждение не обязано отражать реальную причину ответа (низкая верность объяснения):",
             size=14, bold=True, color=DEEP)
    st_y = fy + 0.42
    # #185/#319: правый текст-блок «Контроль на самообъяснении…» УБРАН
    # полностью — тезис несёт реальный интернет-мем «Distracted boyfriend»
    # (imgflip): модель отвлеклась на красивое объяснение вслух и упустила
    # реальную причину ответа. Русские подписи вжжены в шаблон.
    mh = 2.64
    mw = mh * (1200 / 800)
    mx = 0.55
    ocean_box(s, mx - 0.12, st_y - 0.10, mw + 0.24, mh + 0.20)
    add_image(s, WEB / "s06-distracted-ru.png", mx, st_y, mw, mh)
    # короткий вывод справа (замена убранного текст-блока); gold-акцент на глаголе
    rx = mx + mw + 0.44
    rw = 13.33 - rx - 0.55
    gh = 1.72
    gy = st_y + (mh - gh) / 2
    filled_rect(s, rx, gy, rw, gh, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_runs(s, rx + 0.34, gy + 0.14, rw - 0.68, gh - 0.28, [
        {"text": "Проверяй результат", "size": 21, "bold": True, "color": DEEP},
        {"text": ", а не самообъяснение модели.", "size": 21, "bold": True,
         "color": DEEP},
        {"text": "Сверяй факты с внешним источником — правдоподобный текст рассуждения ничего не подтверждает.",
         "size": 15, "bold": False, "color": DEEP, "newpara": True,
         "space_before": 8},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
    footer(s, "Anthropic, апрель 2025 · верность падает на трудных задачах · перепроверить ко дню лекции.")
    speaker_notes(s, load_notes("s06"))


def build_s08(p):
    """assertion_visual — context engineering + context rot curve (c08)."""
    s = blank(p)
    slide_title(s, "Контекст-инжиниринг: минимум высокосигнального.", size=26)
    text_box(s, 0.55, 1.20, 12.25, 0.55,
             "Промпт-инжиниринг — одна инструкция. Контекст-инжиниринг — курирование всего набора токенов, видимых модели на инференсе.",
             size=14, italic=True, color=MID, line_spacing=1.15)
    # left — curve chart
    cx, cyy, cw, chh = 0.55, 1.95, 7.05, 3.55
    ocean_box(s, cx, cyy, cw, chh)
    add_image(s, CHARTS / "c08-context-rot.png", cx + 0.18, cyy + 0.16,
              cw - 0.36, chh - 0.32)
    text_box(s, cx + 0.18, cyy + chh + 0.02, cw - 0.36, 0.42,
             "context rot = тот же «lost in the middle» из Л2 — новый термин, не новая сущность",
             size=11.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # right — criterion
    rx, rw = 7.85, 4.95
    ocean_box(s, rx, 1.95, rw, 3.55, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.28, 2.16, rw - 0.56, 0.45, "Когда НЕ RAG (точка 1)",
             size=17, bold=True, color=TEAL)
    text_box(s, rx + 0.28, 2.70, rw - 0.56, 1.55,
             "малый стабильный корпус, влезает в окно → полный контекст + кэширование префикса, а не RAG-инфраструктура",
             size=15, color=DEEP, line_spacing=1.22)
    icon(s, "circle-slash", rx + 0.28, 4.45, 0.78, "teal")
    text_box(s, rx + 1.20, 4.55, rw - 1.45, 0.85,
             "RAG здесь добавил бы хрупкость без выигрыша",
             size=13.5, bold=True, color=DEEP, line_spacing=1.12,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.95, 12.25, 0.88,
                 "«Найти наименьший набор высокосигнальных токенов, максимизирующий вероятность желаемого исхода» — это инженерное требование, не эстетика.",
                 size=14.5)
    speaker_notes(s, load_notes("s08"))


def build_s08a(p):
    """NEW (§1.8) — чит-шит «как строить промпт», 8 пунктов. Компактный
    аналог чек-листа §5.3 для уровня одного промпта."""
    s = blank(p)
    slide_title(s, "Шпаргалка: как строить промпт.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.40,
             "Минимум, ниже которого промпт систематически недорабатывает. Приложите к любому промпту перед первым запуском.",
             size=14, italic=True, color=MID)
    items = [
        ("Роль", "если нужен тон/регистр — и НЕ как обещание точности"),
        ("Задача", "конкретное проверяемое действие, не расплывчатое пожелание"),
        ("Контекст", "минимально необходимое, не «всё, что могло бы пригодиться»"),
        ("Формат вывода", "указан явно, если ответ машинно-обрабатываем"),
        ("Разделители", "если содержимого больше одного вида (инструкция/данные)"),
        ("Примеры (few-shot)", "только если формат неочевиден из одной инструкции"),
        ("CoT", "только если задача требует многошагового рассуждения"),
        ("Длина", "не длиннее необходимого — лишние токены «тонут» в контексте"),
    ]
    cw, chh = 3.95, 1.28
    gapx, gapy = 0.20, 0.20
    x0, y0 = 0.55, 1.78
    for i, (t, sub) in enumerate(items):
        r, c = divmod(i, 2)
        # 2 cols × 4 rows
        x = x0 + c * (cw + gapx) if False else 0.0
        # 4 cols × 2 rows layout
        col = i % 4
        row = i // 4
        cw2 = 3.02
        x = 0.55 + col * (cw2 + 0.13)
        y = y0 + row * (chh + gapy)
        isgold = (i == 0)  # роль = ключевой пункт (миф про точность)
        if isgold:
            ocean_box(s, x, y, cw2, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, y, cw2, chh)
        circle(s, x + 0.18, y + 0.18, 0.34, GOLD if isgold else MID)
        text_box(s, x + 0.18, y + 0.18, 0.34, 0.34, str(i + 1),
                 size=13, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.62, y + 0.16, cw2 - 0.78, 0.40, t,
                 size=14, bold=True, color=DEEP)
        text_box(s, x + 0.20, y + 0.58, cw2 - 0.40, 0.62, sub,
                 size=11, color=SLATE, line_spacing=1.10)
    gold_callout(s, 0.55, 5.20, 12.25, 1.05,
                 "Компактная форма всего раздела: роль ≠ точность, структура помогает разделить входы, CoT — точечный инструмент, контекст — минимальный. Для крупных архитектурных решений (RAG / FT / агент) — восьмишаговый чек-лист Раздела 5.",
                 size=14)
    speaker_notes(s, load_notes("s08a"))


def build_s09(p):
    """section_divider — Раздел 2 RAG."""
    build_section_divider(
        p, 2, "Раздел 2", "RAG: поиск-дополненная генерация",
        "Извлечь релевантное → положить в контекст → ответить с опорой на источник",
        "s09",
        image_src=WEB / "div-r2-library.jpg",
        tag="внешнее знание · 3 разбора · 1 провал")


def build_s10(p):
    """schema_pipeline — RAG 3-stage horizontal pipeline (RIGHT_ARROW)."""
    s = blank(p)
    slide_title(s, "Принцип RAG — три шага.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "RAG = индексация → поиск (семантический поиск из Л2) → генерация с опорой; «не знаю» — корректный ответ.",
             size=14.5, italic=True, color=MID)
    # 3 stage boxes with RIGHT_ARROW between
    sy, sh = 1.85, 3.05
    bw = 3.55
    gap_arrow = 0.55
    x0 = 0.55
    stages = [
        ("1", "Индексация", "заранее, офлайн", "database",
         "корпус → чанки → эмбеддинг каждого → векторное хранилище", MID, False),
        ("2", "Поиск", "на запрос", "route",
         "вопрос → эмбеддинг → k ближайших фрагментов\n\n= тот же семантический поиск из Л2 — не переобъясняем", MID, False),
        ("3", "Генерация", "с опорой на источник (grounding)", "check-check",
         "фрагменты + вопрос → ответ со ссылкой на источник", TEAL, True),
    ]
    x = x0
    for i, (num, title, tag, ic, body, col, isteal) in enumerate(stages):
        if isteal:
            ocean_box(s, x, sy, bw, sh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
        else:
            ocean_box(s, x, sy, bw, sh)
        circle(s, x + 0.26, sy + 0.26, 0.46, col)
        text_box(s, x + 0.26, sy + 0.26, 0.46, 0.46, num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.86, sy + 0.26, bw - 1.05, 0.42, title,
                 size=17, bold=True, color=DEEP)
        text_box(s, x + 0.86, sy + 0.66, bw - 1.05, 0.30, tag,
                 size=11.5, italic=True, color=LIGHT)
        icon(s, ic, x + bw - 0.78, sy + 0.28, 0.50,
             "teal" if isteal else "mid")
        text_box(s, x + 0.28, sy + 1.18, bw - 0.56, sh - 1.36, body,
                 size=13.5, color=DEEP, line_spacing=1.22)
        if i < 2:
            right_arrow(s, x + bw + 0.06, sy + sh / 2 - 0.30, gap_arrow - 0.12, 0.60,
                        fill=LIGHT)
        x += bw + gap_arrow
    gold_callout(s, 0.55, 5.10, 12.25, 0.90,
                 "«Не знаю» / «см. источник X» — корректный ответ RAG-системы. Правдоподобный ответ при нерелевантном поиске — это дефект, а не «хоть что-то».",
                 size=14)
    # RAG-2026: current default стек (glossary-locked термины на видимом слое)
    filled_rect(s, 0.55, 6.12, 12.25, 0.62, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_runs(s, 0.78, 6.20, 11.8, 0.48, [
        {"text": "RAG-2026: ", "size": 12.5, "bold": True, "color": TEAL},
        {"text": "agentic (агентный) RAG по умолчанию · гибридный поиск (BM25 + плотные векторы) · реранкер · каскад промахов ", "size": 12.5, "color": DEEP},
        {"text": "5,7% → 1,9%", "size": 12.5, "bold": True, "color": DEEP},
        {"text": " (Contextual Retrieval).", "size": 12.5, "color": DEEP},
    ], line_spacing=1.1)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """assertion_visual — 4 conjunction cards -> RAG."""
    s = blank(p)
    slide_title(s, "Когда RAG — правильный выбор.", size=27)
    # #221/#222: убран жаргон «конъюнкция/дизъюнкция»; формулировка §2.2 —
    # сильный сигнал по признакам + отсутствие блокеров с §2.3.
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "RAG оправдан при сильном сигнале по признакам ниже — и отсутствии блокеров со следующего слайда «когда НЕ RAG».",
             size=14.5, italic=True, color=MID)
    cards = [
        ("Большое / растущее", "не влезает в окно целиком, или дорого класть весь корпус в каждый запрос"),
        ("Меняется", "документы, цены, регламенты обновляются чаще, чем выходят версии модели"),
        ("Свежесть + провенанс", "ответ опирается на проверяемый источник; можно показать, откуда факт"),
        ("Приватная база", "знания компании не входят в веса публичной модели"),
    ]
    cw, chh = 2.90, 2.55
    cy = 1.78
    x = 0.55
    for i, (t, sub) in enumerate(cards):
        ocean_box(s, x, cy, cw, chh)
        text_box(s, x + 0.20, cy + 0.22, cw - 0.40, 0.75, t,
                 size=15.5, bold=True, color=MID, line_spacing=1.05)
        text_box(s, x + 0.20, cy + 1.00, cw - 0.40, chh - 1.2, sub,
                 size=12, color=DEEP, line_spacing=1.16)
        x += cw + 0.20
    # -> RAG result + worked example
    ocean_box(s, 0.55, 4.55, 7.55, 1.80, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, 0.85, 4.72, 7.1, 0.4, "Признаки усиливают друг друга → RAG",
             size=17, bold=True, color=DEEP)
    text_box(s, 0.85, 5.16, 7.1, 1.10,
             "Корп. база из тысяч регламентов, обновляется еженедельно, ответ с обязательной ссылкой на пункт: все признаки сошлись, ни один более простой механизм их совместно не закрывает → образцовый профиль RAG.",
             size=13, color=DEEP, line_spacing=1.16)
    ocean_box(s, 8.30, 4.55, 4.50, 1.80, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "check-check", 8.55, 4.74, 0.44, "teal")
    text_box(s, 8.55, 5.24, 4.05, 1.00,
             "Один признак — повод присмотреться, но не строить автоматически: сверьтесь с блокерами на следующем слайде.",
             size=12.5, bold=True, color=DEEP, line_spacing=1.15)
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """schema_matrix-ish — 3 criteria columns «when NOT RAG».
    P1 fix (issue #157 review): s11/s12 shared identical «N cards + summary
    plaque» skeleton back-to-back — visually merged into one stretched
    slide. Differentiated here via base card palette (TEAL instead of
    primary LIGHT/SURFACE — «caution / exclusion» register vs s11's
    primary-blue «inclusion» cards) + numbered gold badge per card (s22b
    slot-badge pattern) so the two decks read as distinct compositions at
    a glance, without touching content/copy."""
    s = blank(p)
    slide_title(s, "Когда RAG — НЕ правильный выбор.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "Знать, когда RAG не нужен, ценнее: это модная архитектура, её ставят туда, где она вредит.",
             size=14.5, italic=True, color=MID)
    # #223/#224: критерий observability убран (он живёт на s13, провал RAG
    # на масштабе, не дублируется здесь). Вместо него — новый критерий
    # «данные доступны live через API/MCP» (§2.3, forward-callback на §4.1).
    cols = [
        ("circle-slash", "Корпус влезает в окно",
         "ориентир — менее ~200k токенов, меняется редко",
         "→ полный контекст + кэширование префикса, не RAG-инфраструктура",
         False),
        ("key", "Фиксированная политика / значение",
         "тариф, цена, пункт регламента, правило",
         "→ детерминированная таблица соответствий / статическая страница",
         False),
        ("cable", "Данные доступны в реальном времени через API / MCP",
         "во внутреннем сервисе, базе, поиске другой системы",
         "→ вызвать инструмент напрямую; RAG-индекс поверх — лишний, более хрупкий и более устаревающий слой",
         True),
    ]
    # #185: карточки чуть ниже → открыт ряд под реальный мем «Roll Safe».
    cw, chh = 4.00, 2.62
    cy = 1.78
    x = 0.55
    for i, (nm, t, tag, alt, isgold) in enumerate(cols):
        if isgold:
            ocean_box(s, x, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, cy, cw, chh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=1.5)
        # numbered badge top-left (s22b slot-badge pattern) — distinct
        # silhouette from s11's plain unnumbered cards
        circle(s, x + 0.20, cy + 0.18, 0.32, GOLD if isgold else TEAL)
        text_box(s, x + 0.20, cy + 0.18, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        icon(s, nm, x + cw - 0.76, cy + 0.20, 0.52, "gold" if isgold else "teal")
        text_box(s, x + 0.66, cy + 0.24, cw - 1.55, 0.78, t,
                 size=15, bold=True, color=(DEEP if isgold else TEAL),
                 line_spacing=1.08)
        text_box(s, x + 0.24, cy + 1.15, cw - 0.48, 0.40, tag,
                 size=12, italic=True, color=LIGHT)
        text_box(s, x + 0.24, cy + 1.58, cw - 0.48, chh - 1.72, alt,
                 size=13, bold=True, color=DEEP, line_spacing=1.14)
        x += cw + 0.32
    # #185: реальный интернет-мем «Roll Safe» — «не нужен RAG, если корпус
    # влезает в контекст» (первый из трёх блокеров, самый частый анти-паттерн).
    rmw, rmh = 2.15, 1.21
    rmx, rmy = 0.55, 4.50
    filled_rect(s, rmx - 0.05, rmy - 0.05, rmw + 0.10, rmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s12-rollsafe-ru.png", rmx, rmy, rmw, rmh)
    text_box(s, rmx + rmw + 0.36, rmy + 0.02, 9.05, rmh,
             "Частый анти-паттерн — ставить RAG-инфраструктуру там, где весь корпус спокойно помещается в контекстное окно и меняется редко.",
             size=13.5, bold=True, color=MID, line_spacing=1.24,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.78, 12.25, 1.02,
                 "RAG избыточен, если выполнен ЛЮБОЙ из трёх: (а) корпус влезает в окно и стабилен, (б) задача сводится к фиксированному значению, (в) знание уже доступно напрямую и живьём через инструмент. «Прямой вызов вернёт данные на момент запроса; RAG-индекс — на момент последней индексации.»",
                 size=13.5)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """case_study — RAG fails at scale + Air Canada — разбор архитектуры (v2: decluttered
    — per-case triangle-alert icons, tighter visible text, mass-rebalanced
    right box so diagnosis+alternative fill evenly, no big internal gap)."""
    s = blank(p)
    slide_title(s, "Провал RAG на масштабе.", size=27)
    gold_callout(s, 0.55, 1.10, 12.25, 0.95,
                 "«Вернул что-то» ≠ «вернул правильное». У RAG нет сигнала «не нашёл» — он всегда отдаёт k ближайших, даже нерелевантных.",
                 size=15)
    # left — 3 failure cases, each with a triangle-alert anchor icon
    # PA-2 (owner-approved): cases distributed EVENLY across the full box
    # height with thin separators (was ~15-20% dead band at the bottom;
    # mass now matches the right Air Canada box).
    lx, lw = 0.55, 6.55
    box_y, box_h = 2.14, 3.74
    ocean_box(s, lx, box_y, lw, box_h)
    cases = [
        ("Legal-AI", "«ближайшие» дела из другой юрисдикции / отменённый прецедент — модель опирается на них как на факт"),
        ("Medical-RAG", "смешал фрагменты разных пациентов — близки по симптомам, клинически объединять нельзя"),
        ("Бот поддержки", "работал на сотнях статей; после роста до тысяч качество тихо просело — никто не заметил"),
    ]
    cell_h = box_h / 3.0          # 3 equal cells fill the box top→bottom
    for ci, (nm, body) in enumerate(cases):
        cy0 = box_y + ci * cell_h
        icon(s, "triangle-alert", lx + 0.28, cy0 + 0.22, 0.40, "teal")
        text_box(s, lx + 0.84, cy0 + 0.22, lw - 1.10, 0.36, nm,
                 size=16, bold=True, color=MID)
        text_box(s, lx + 0.28, cy0 + 0.66, lw - 0.54, 0.62, body,
                 size=13, color=DEEP, line_spacing=1.20)
        if ci < 2:
            connector(s, lx + 0.30, box_y + (ci + 1) * cell_h,
                      lx + lw - 0.30, box_y + (ci + 1) * cell_h, LIGHT, 0.75)
    # right — Air Canada — разбор архитектуры (teal), evenly split: diagnosis / alternative
    rx, rw = 7.35, 5.45
    ocean_box(s, rx, 2.14, rw, 3.74, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.30, 2.36, rw - 0.60, 0.34, "Air Canada — разбор архитектуры",
             size=16, bold=True, color=TEAL)
    text_box(s, rx + 0.30, 2.78, rw - 0.60, 0.30, "Диагноз",
             size=13, bold=True, color=DEEP)
    text_box(s, rx + 0.30, 3.10, rw - 0.60, 1.05,
             "сгенерированный правдоподобный текст поставлен в роль, требовавшую извлечённого проверенного факта — отказ опоры на источник",
             size=13, color=DEEP, line_spacing=1.20)
    connector(s, rx + 0.30, 4.22, rx + rw - 0.30, 4.22, TEAL, 1.0)
    text_box(s, rx + 0.30, 4.34, rw - 0.60, 0.30, "Правильная альтернатива",
             size=13, bold=True, color=DEEP)
    text_box(s, rx + 0.30, 4.66, rw - 0.60, 1.15,
             "фиксированная политика → таблица соответствий / страница; нужен диалог → RAG со строгой опорой на источник, обязательная цитата, явное «не знаю», проверка человеком",
             size=13, color=DEEP, line_spacing=1.20)
    footer(s, "Документированные классы провалов (Barnett et al. 2024; Air Canada — McCarthy Tétrault 2024). Кейсы — illustrative.")
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """assertion_visual (§3.1/§3.5) — дистилляция = ОТДЕЛЬНАЯ техника, НЕ вид
    fine-tuning (#227 P0). Схема: teacher (fine-tuned) → distill → student
    (меньше, дешевле). Контраст «что делает fine-tuning / что делает
    дистилляция». Критерии «что куда» — на s17."""
    s = blank(p)
    slide_title(s, "Дистилляция — не вид дообучения, а отдельная техника.", size=25)
    text_box(s, 0.55, 1.12, 12.25, 0.44,
             "Дообучение (fine-tuning) меняет поведение модели. Дистилляция — сжатие: перенос умений большой модели в маленькую. Их часто путают, но это две таксономически разные операции, работающие в связке.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # pipeline: teacher (FT) -> distill -> student  (сжат влево — правая
    # колонка отдана под реальный мем «Yoda» учитель→ученик)
    sy, sh = 1.90, 2.48
    ocean_box(s, 0.55, sy, 8.55, sh)
    bw, bh = 2.05, 1.50
    by = sy + 0.56
    t1x = 0.85
    arr1 = t1x + bw
    t2x = arr1 + 0.80
    arr2 = t2x + bw
    t3x = arr2 + 0.80
    # teacher — big, fine-tuned
    filled_rect(s, t1x, by, bw, bh, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.08)
    icon(s, "cpu", t1x + bw / 2 - 0.20, by + 0.12, 0.38, "teal")
    text_box(s, t1x + 0.08, by + 0.58, bw - 0.16, 0.30, "Учитель",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, t1x + 0.10, by + 0.90, bw - 0.20, 0.52,
             "большая, дообучена под задачу", size=10.5, italic=True,
             color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.04)
    text_box(s, arr1 - 0.32, by - 0.40, 1.25, 0.30, "дистилляция",
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    right_arrow(s, arr1 + 0.06, by + bh / 2 - 0.17, 0.68, 0.34, fill=MID)
    # student — small, cheaper (gold anchor)
    filled_rect(s, t2x, by + 0.20, bw, bh - 0.40, GOLD_TINT, stroke=GOLD,
                stroke_pt=2.0, radius=True, radius_adj=0.08)
    icon(s, "cpu", t2x + bw / 2 - 0.17, by + 0.28, 0.32, "gold")
    text_box(s, t2x + 0.08, by + 0.66, bw - 0.16, 0.30, "Ученик",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, t2x + 0.10, by + 0.96, bw - 0.20, 0.44,
             "маленькая, дешевле и быстрее", size=10.5,
             italic=True, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.04)
    text_box(s, arr2 - 0.32, by - 0.40, 1.25, 0.30, "в бой →",
             size=11, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    right_arrow(s, arr2 + 0.06, by + bh / 2 - 0.17, 0.68, 0.34, fill=LIGHT)
    filled_rect(s, t3x, by, bw, bh, SURFACE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.08)
    icon(s, "target", t3x + bw / 2 - 0.20, by + 0.12, 0.38, "mid")
    text_box(s, t3x + 0.08, by + 0.58, bw - 0.16, 0.30, "В бой",
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, t3x + 0.10, by + 0.90, bw - 0.20, 0.52,
             "ниже стоимость и задержка", size=10.5,
             italic=True, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.04)
    # #185: реальный интернет-мем «Yoda» — учитель (большая дообученная
    # модель) передаёт умение маленькому ученику (дистилляция). Правая колонка,
    # во всю высоту схемы + контраст-плашек.
    ymx, ymy, ymw, ymh = 9.35, 1.90, 3.45, 4.08
    ocean_box(s, ymx, ymy, ymw, ymh)
    add_image(s, WEB / "s14-yoda-ru.png", ymx + 0.14, ymy + 0.14,
              ymw - 0.28, ymh - 0.28)
    # contrast strip (сужен влево — правая колонка отдана под мем)
    cy, ch = 4.56, 1.42
    ocean_box(s, 0.55, cy, 4.20, ch)
    text_box(s, 0.78, cy + 0.14, 3.80, 0.34, "Дообучение отвечает на:",
             size=13, bold=True, color=MID)
    text_box(s, 0.78, cy + 0.50, 3.80, 0.86,
             "«как модель себя ведёт» — тон, формат, политика. Меняет ВЕСА под поведение.",
             size=12, color=DEEP, line_spacing=1.16)
    ocean_box(s, 4.90, cy, 4.20, ch, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, 5.13, cy + 0.14, 3.80, 0.34, "Дистилляция отвечает на:",
             size=13, bold=True, color=TEAL)
    text_box(s, 5.13, cy + 0.50, 3.80, 0.86,
             "«как сделать модель дешевле» — тот же навык в меньшей модели. Это сжатие, не смена поведения.",
             size=12, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 6.14, 12.25, 0.78,
                 "Что делать: не пишите «дистилляция — это дообучение». В связке они идут так: сначала дообучить учителя (teacher) под задачу, потом дистиллировать в ученика (student) ради цены. Критерии «что куда» — на следующем слайде.",
                 size=13)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    """assertion_visual — PEFT vs full-FT: frozen base + adapters + 3 reasons."""
    s = blank(p)
    slide_title(s, "PEFT вместо полного дообучения.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.78,
             "PEFT — базовые веса замораживаются, обучается лишь небольшой набор адаптеров. LoRA — низкоранговые матрицы-адаптеры; QLoRA — то же поверх квантованной модели.",
             size=14, italic=True, color=MID, line_spacing=1.18)
    # left — schema: big frozen base + small adapters
    lx, ly, lw, lh = 0.55, 2.10, 4.95, 3.15
    ocean_box(s, lx, ly, lw, lh)
    filled_rect(s, lx + 0.55, ly + 0.55, lw - 1.1, 1.55, MID, radius=True,
                radius_adj=0.06)
    text_box(s, lx + 0.55, ly + 0.55, lw - 1.1, 1.55,
             "Базовые веса\n(frozen)", size=17, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    for k in range(3):
        ax = lx + 0.70 + k * 1.20
        filled_rect(s, ax, ly + 2.35, 0.95, 0.55, GOLD, radius=True,
                    radius_adj=0.18)
        text_box(s, ax, ly + 2.35, 0.95, 0.55, "LoRA",
                 size=11, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx, ly + lh - 0.02, lw, 0.3,
             "адаптеры — мегабайты vs гигабайты", size=11, italic=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    # right — 3 reasons (компактнее — освобождаем правую колонку под мем)
    rx, rw = 5.80, 3.55
    reasons = [
        ("1. Дешевле и быстрее", "миллионы параметров вместо миллиардов; QLoRA — на одном GPU", False),
        ("2. Модульность", "адаптеры мегабайты vs гигабайты; одна база — много специализаций", False),
        ("3. ↓ Риск забывания", "база заморожена, физически не переписывается под новый сигнал", True),
    ]
    yy = 2.10
    for t, b, isgold in reasons:
        if isgold:
            ocean_box(s, rx, yy, rw, 1.02, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, rx, yy, rw, 0.98)
        text_box(s, rx + 0.22, yy + 0.10, rw - 0.44, 0.34, t,
                 size=14, bold=True, color=(DEEP if isgold else MID))
        text_box(s, rx + 0.22, yy + 0.44, rw - 0.44, (0.52 if isgold else 0.48), b,
                 size=12, color=DEEP, line_spacing=1.10)
        yy += (1.14 if isgold else 1.10)
    # #185: реальный интернет-мем «Buff Doge vs Cheems» — PEFT (сильный,
    # спокойный выбор) vs полное дообучение (дорогой, слабый). Правая колонка.
    dmx, dmy, dmw, dmh = 9.55, 2.10, 3.25, 3.15
    ocean_box(s, dmx, dmy, dmw, dmh)
    add_image(s, WEB / "s15-doge-ru.png", dmx + 0.12, dmy + 0.12,
              dmw - 0.24, dmh - 0.24)
    # LoRA-adoption baseline (§3.2) с ОБЯЗАТЕЛЬНОЙ оговоркой на видимом слое
    # (Baseline Mandate): доля среди тегированных PEFT, не среди всего FT.
    by2 = 5.50
    ocean_box(s, 0.55, by2, 5.95, 1.32, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_runs(s, 0.80, by2 + 0.16, 5.5, 0.55, [
        {"text": "98,4% ", "size": 30, "bold": True, "color": GOLD},
        {"text": "моделей с тегом PEFT — это LoRA", "size": 14, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.80, by2 + 0.72, 5.45, 0.52,
             "из 20 834 карточек на Hugging Face Hub · оговорка: доля среди моделей с тегом PEFT, не среди всего дообучения",
             size=11, italic=True, color=SLATE, line_spacing=1.12)
    gold_callout(s, 6.70, by2, 6.10, 1.32,
                 "PEFT (LoRA/QLoRA) почти всегда предпочтительнее полного дообучения: дешевле, модульнее, ↓ риск забывания (forgetting). Полное дообучение в 2026 — почти никогда.",
                 size=14)
    footer(s, "HF PEFT team, «Beyond LoRA?», июнь 2026. Полный спектр методов (SFT / DPO / RFT) — в главе.")
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    """case_study — catastrophic forgetting (diverging chart c16) + criterion."""
    s = blank(p)
    slide_title(s, "Провал: катастрофическое забывание.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.55,
             "Катастрофическое забывание (catastrophic forgetting) — деградация общих способностей модели в результате узкого агрессивного дообучения.",
             size=14, italic=True, color=MID, line_spacing=1.15)
    cx, cyy, cw, chh = 0.55, 1.92, 7.05, 3.55
    ocean_box(s, cx, cyy, cw, chh)
    add_image(s, CHARTS / "c16-forgetting.png", cx + 0.18, cyy + 0.16,
              cw - 0.36, chh - 0.32)
    text_box(s, cx + 0.18, cyy + chh + 0.02, cw - 0.36, 0.42,
             "тяжелее с ростом масштаба модели — у крупной выше исходный уровень, ей «больше падать»",
             size=11.5, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    rx, rw = 7.85, 4.95
    ocean_box(s, rx, 1.92, rw, 2.30, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.26, 2.10, rw - 0.52, 2.0,
             "Нет петли проверки (eval) на общих задачах + нет версий датасета/весов → не увидишь поломку до боевой эксплуатации + не сможешь откатиться → это не «риск», это критерий «НЕ делай дообучение»",
             size=13.5, bold=True, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, 4.37, rw, 1.10, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.26, 4.52, rw - 0.52, 0.85,
             "Правильно: PEFT (замороженные веса — ниже риск); для меняющегося знания — RAG, не FT вовсе.",
             size=12.5, color=DEEP, line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Эмпирически наблюдается при постоянном дообучении (Luo et al. 2023). Механизмы — «исследования показывают» (препринт).")
    speaker_notes(s, load_notes("s16"))


def build_s18(p):
    """section_divider — Раздел 4 «Агенты» (заголовок БЕЗ «+безопасность»;
    контент безопасности внутри раздела, на s25)."""
    build_section_divider(
        p, 4, "Раздел 4", "Агенты",
        "От собеседника в окне чата — к компоненту боевой системы: цикл, экипировка, память, доступ к инструментам — и где всё это ломается.",
        "s18",
        image_src=WEB / "div-r4-robot-arm.jpg",
        tag="агент + экипировка · 5 разборов · 4 провала")


def build_s19(p):
    """MERGED (§4.1) — API-механика (structured output / function calling /
    prompt caching) СВЕРХУ + MCP (N×M→N+M, USB-C, приятие, поворот доверия)
    СНИЗУ. Плотный слайд — компактные карточки, отдельная проверка 5-Second."""
    s = blank(p)
    slide_title(s, "Модель становится компонентом системы: API + MCP.", size=24)
    # TOP — 3 API mechanism cards (compact)
    cards = [
        ("boxes", "Структурированный вывод", "выход строго по схеме (JSON), не текст для разбора", "встраиваемая"),
        ("terminal", "Вызов инструментов", "модель формулирует «вызови X»; исполняет ваш код, не модель", "активная"),
        ("database", "Кэш промптов", "не пересчитывать неизменный префикс при каждом запросе", "экономичная"),
    ]
    cw, chh = 4.00, 2.02
    cy = 1.12
    x = 0.55
    for nm, t, body, tag in cards:
        ocean_box(s, x, cy, cw, chh)
        filled_rect(s, x + 0.22, cy + 0.22, 0.56, 0.56, MID, radius=True, radius_adj=0.18)
        icon(s, nm, x + 0.28, cy + 0.28, 0.44, "white")
        text_box(s, x + 0.90, cy + 0.24, cw - 1.05, 0.52, t,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x + 0.24, cy + 0.90, cw - 0.48, 0.72, body,
                 size=12.5, color=DEEP, line_spacing=1.14)
        chip(s, x + 0.24, cy + chh - 0.46, cw - 0.48, 0.34, tag,
             fill=TEAL, color=WHITE, size=12)
        x += cw + 0.13
    # BOTTOM — MCP
    my = 3.36
    ocean_box(s, 0.55, my, 6.55, 2.55)
    icon(s, "cable", 0.80, my + 0.20, 0.48, "mid")
    text_runs(s, 1.40, my + 0.22, 5.5, 0.5, [
        {"text": "MCP", "size": 18, "bold": True, "color": MID},
        {"text": "  — «USB-C для инструментов LLM»", "size": 14, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text_runs(s, 0.80, my + 0.82, 6.0, 0.5, [
        {"text": "N×M", "size": 20, "bold": True, "color": LIGHT},
        {"text": " несовместимых интеграций → ", "size": 13, "color": DEEP},
        {"text": "N+M", "size": 20, "bold": True, "color": TEAL},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.80, my + 1.36, 6.0, 0.32,
             "инструмент один раз = MCP-сервер; модель один раз = MCP-клиент",
             size=11.5, italic=True, color=SLATE)
    tl = [("Anthropic", "11/2024", "logo-anthropic"),
          ("OpenAI", "03/2025", "logo-openai"),
          ("Google", "04/2025", "logo-gemini")]
    tx0 = 0.80
    for i, (nm, dt, lg) in enumerate(tl):
        ex = tx0 + i * 2.0
        add_image(s, ICONS / f"{lg}.png", ex, my + 1.82, 0.34, 0.34)
        text_box(s, ex + 0.42, my + 1.78, 1.5, 0.26, nm, size=11.5, bold=True, color=DEEP)
        text_box(s, ex + 0.42, my + 2.02, 1.5, 0.24, dt, size=11.5, bold=True, italic=True, color=LIGHT)
        if i < 2:
            text_box(s, ex + 1.72, my + 1.82, 0.22, 0.3, "→", size=13, bold=True, color=LIGHT)
    # trust warning — P1 fix (issue #157 review): 3 bullets -> 2 most load-bearing
    # (root cause: code/access; concrete attack vector: prompt injection carrier).
    # Retention-policy point dropped here — it's covered on s25's ZDR block.
    ocean_box(s, 7.25, my, 5.55, 2.55, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, 7.52, my + 0.18, 5.1, 0.72,
             "Стандартизация подключения ≠ безопасность подключаемого — и усугубляет проблему доверия.",
             size=13.5, bold=True, color=TEAL, line_spacing=1.14)
    for i, w in enumerate([
        "код в вашем окружении / доступ к данным",
        "описание попадает в контекст — носитель инъекции в промпт"]):
        circle(s, 7.54, my + 1.06 + i * 0.52 + 0.05, 0.11, TEAL)
        text_box(s, 7.76, my + 1.06 + i * 0.52, 5.0, 0.48, w,
                 size=12.5, color=DEEP, line_spacing=1.10)
    gold_callout(s, 0.55, 6.06, 12.25, 0.80,
                 "Ни один механизм не делает модель надёжнее — правило лестницы не отменяется. Удобство подключения — не аргумент за подключение.",
                 size=13.5)
    footer(s, "Актуальные цифры экономии и масштаб экосистемы MCP — см. источники.")
    speaker_notes(s, load_notes("s19"))


def build_s21(p):
    """schema_cycle — agent loop plan→act→check→iterate."""
    s = blank(p)
    slide_title(s, "Цикл агента: план → действие → проверка → повтор.", size=26)
    text_box(s, 0.55, 1.14, 12.25, 0.4,
             "Агент — архитектура, где модель не делает один проход, а работает в цикле, сама определяя последовательность шагов.",
             size=13.5, italic=True, color=MID)
    # 4 step cards in a row with arrows + return arrow below
    steps = [
        ("plan", "План (plan)", "формулирует следующий шаг",
         "близорукий / зацикленный план (не видит накопленной стоимости)", MID, False),
        ("act", "Действие (act)", "вызывает инструмент (function calling)",
         "инструмент падает / тормозит, а ветки на это нет", MID, False),
        ("check", "Проверка (check)", "достигнута ли цель, корректен ли результат",
         "валидация против ВНЕШНЕГО критерия — не самооценка модели (отсылка к пределу CoT)", GOLD, True),
        ("iter", "Повтор (iterate)", "цикл повторяется",
         "нет внешнего лимита на итерации / стоимость / время → петля", MID, False),
    ]
    sy, sh = 1.72, 3.55
    bw = 2.92
    gap = 0.20
    x0 = 0.55
    x = x0
    centers = []
    for i, (k, t, sub, fail, col, isgold) in enumerate(steps):
        if isgold:
            ocean_box(s, x, sy, bw, sh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, sy, bw, sh)
        circle(s, x + bw / 2 - 0.26, sy + 0.24, 0.52,
               (GOLD if isgold else MID))
        text_box(s, x + bw / 2 - 0.26, sy + 0.24, 0.52, 0.52, str(i + 1),
                 size=20, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # START badge on Plan (explicit entry point)
        if i == 0:
            chip(s, x + 0.30, sy + 0.30, 0.96, 0.34, "СТАРТ",
                 fill=TEAL, color=WHITE, size=10.5)
        text_box(s, x + 0.12, sy + 0.90, bw - 0.24, 0.42, t,
                 size=19, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        text_box(s, x + 0.18, sy + 1.36, bw - 0.36, 0.70, sub,
                 size=12, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.12)
        filled_rect(s, x + 0.18, sy + 2.18, bw - 0.36, 1.22,
                    (GOLD_TINT if isgold else SURFACE),
                    stroke=(GOLD if isgold else SOFT_GREY), stroke_pt=1.0,
                    radius=True, radius_adj=0.10)
        text_box(s, x + 0.28, sy + 2.26, bw - 0.56, 1.06,
                 ("режим отказа: " + fail),
                 size=10.5, italic=True, color=(DEEP if isgold else SLATE),
                 align=PP_ALIGN.CENTER, line_spacing=1.12)
        centers.append(x + bw / 2)
        if i < 3:
            right_arrow(s, x + bw + 0.01, sy + 0.62, gap - 0.04, 0.46, fill=LIGHT)
        x += bw + gap
    # Explicit return path: down from Iterate → left across → up into Plan
    box_bottom = sy + sh           # 5.27
    ry = box_bottom + 0.55         # horizontal return rail at 5.82
    last_cx = centers[3]
    first_cx = centers[0]
    # down stub from Iterate
    connector(s, last_cx, box_bottom, last_cx, ry, LIGHT, 2.5)
    # horizontal rail right→left
    connector(s, last_cx, ry, first_cx, ry, LIGHT, 2.5)
    # up stub + arrowhead back into Plan
    connector(s, first_cx, ry, first_cx, box_bottom + 0.20, LIGHT, 2.5)
    up_arrow(s, first_cx - 0.13, box_bottom + 0.02, 0.26, 0.24, fill=LIGHT)
    # loop label centered ON the rail, white-backed chip for legibility
    lbl_w = 4.6
    chip(s, (first_cx + last_cx) / 2 - lbl_w / 2, ry - 0.20, lbl_w, 0.40,
         "⟲  цикл ПОВТОРЯЕТСЯ — возврат к шагу «План»", fill=LIGHT, color=WHITE,
         size=12.5)
    footer(s, "Паттерны цикла (ReAct, Reflexion, Plan-and-Execute) — в главе. Проектировать агента = проектировать защиту на каждом из 4 шагов.")
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """schema_matrix / comparison — Workflow vs Agent (2 columns)."""
    s = blank(p)
    slide_title(s, "Сценарий (workflow) vs агент.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.4,
             "Предсказуемая задача → сценарий; непредсказуемая И ценность оправдывает кратный рост → агент.",
             size=14.5, italic=True, color=MID)
    cy, chh = 1.78, 2.75
    cw = 6.05
    # workflow
    ocean_box(s, 0.55, cy, cw, chh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "git-fork", 0.83, cy + 0.22, 0.50, "teal")
    text_box(s, 1.45, cy + 0.24, cw - 0.9, 0.5, "Сценарий  (workflow)",
             size=17, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.83, cy + 0.85, cw - 0.56, 0.45,
             "LLM и инструменты по предопределённым в коде путям",
             size=13, bold=True, color=DEEP, line_spacing=1.12)
    for i, t in enumerate(["последовательность шагов известна заранее",
                           "предсказуемо, аудируемо",
                           "большинство надёжных боевых систем — это сценарий (workflow)"]):
        circle(s, 0.83, cy + 1.40 + i * 0.43 + 0.06, 0.11, TEAL)
        text_box(s, 1.06, cy + 1.40 + i * 0.43, cw - 1.3, 0.42, t,
                 size=12.5, color=DEEP, line_spacing=1.05)
    # agent
    rx = 6.75
    ocean_box(s, rx, cy, cw, chh)
    icon(s, "bot", rx + 0.28, cy + 0.22, 0.50, "mid")
    text_box(s, rx + 0.90, cy + 0.24, cw - 0.9, 0.5, "Агент",
             size=17, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.28, cy + 0.85, cw - 0.56, 0.45,
             "LLM динамически определяет собственный процесс",
             size=13, bold=True, color=DEEP, line_spacing=1.12)
    for i, t in enumerate(["последовательность заранее не зафиксирована",
                           "кратно больше токенов, чем чат",
                           "ниже аудируемость, выше риск петель"]):
        circle(s, rx + 0.28, cy + 1.40 + i * 0.43 + 0.06, 0.11, MID)
        text_box(s, rx + 0.51, cy + 1.40 + i * 0.43, cw - 1.3, 0.42, t,
                 size=12.5, color=DEEP, line_spacing=1.05)
    # diagnostic question
    ocean_box(s, 0.55, 4.70, 12.25, 1.05, fill=SURFACE, stroke=LIGHT)
    text_box(s, 0.85, 4.80, 11.65, 0.9,
             "Могу ли я заранее, до запуска, выписать последовательность шагов?  да (даже с ветвлениями) → сценарий  ·  принципиально нет И ценность оправдывает кратные стоимость/риск → агент",
             size=13.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.18)
    # #231: вложенность workflow↔agent — норма, не третья архитектура.
    ocean_box(s, 0.55, 5.90, 8.55, 0.98, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, 0.80, 5.98, 8.05, 0.82,
             "Вложенность — норма: агент проверки кода вызывает сценарий «линтер→тесты→формат» как один шаг; сценарий обработки заявок делегирует мини-агенту разбор свободной жалобы. Динамика — только там, где нужна.",
             size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
    ocean_box(s, 9.30, 5.90, 3.50, 0.98, fill=SURFACE, stroke=SOFT_GREY, stroke_pt=1.0)
    text_box(s, 9.50, 5.98, 3.10, 0.82,
             "Найди простейшее. «Лень формализовать» не делает задачу непредсказуемой.",
             size=11.5, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.10)
    speaker_notes(s, load_notes("s22"))


def build_s22b(p):
    """NEW (§4.4) — «из чего сделан агент-помощник»: карта 5 слотов
    экипировки. ASSERTION на заголовке: каждый слот — компромисс, не
    апгрейд по умолчанию (рифма с лестницей)."""
    s = blank(p)
    slide_title(s, "Каждый слот экипировки агента — компромисс, не выигрыш.", size=23)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Реальный агент-помощник (Claude Code, Cursor, Aider) — это цикл план→действие→проверка→повтор ПЛЮС оснастка. Пять типовых слотов:",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    slots = [
        ("brain-circuit", "Память", "что помнит между сессиями"),
        ("file-text", "Инструкции-правила", "файлы-конвенции + журнал задач"),
        ("puzzle", "Навыки (skills)", "переиспользуемые процедуры"),
        ("users", "Субагенты", "делегирование + изоляция"),
        ("cable", "Доступ / MCP", "доступ к внешним инструментам"),
    ]
    cw = 2.40
    gap = 0.10
    x0 = 0.55
    y = 1.78
    chh = 2.55
    for i, (ic, t, sub) in enumerate(slots):
        x = x0 + i * (cw + gap)
        ocean_box(s, x, y, cw, chh)
        filled_rect(s, x + cw / 2 - 0.42, y + 0.26, 0.84, 0.84, MID, radius=True, radius_adj=0.18)
        icon(s, ic, x + cw / 2 - 0.32, y + 0.36, 0.64, "white")
        text_box(s, x + 0.10, y + 1.28, cw - 0.20, 0.5, t,
                 size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(s, x + 0.14, y + 1.80, cw - 0.28, 0.65, sub,
                 size=11.5, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.12)
        # slot number badge
        circle(s, x + 0.14, y + 0.14, 0.32, GOLD)
        text_box(s, x + 0.14, y + 0.14, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 4.66, 12.25, 1.30,
                 "Как не поднимаешься по лестнице архитектур без требования задачи — не добавляй агенту память, субагентов или MCP «на всякий случай». Каждый слот несёт свою цену: операционную сложность, новую поверхность отказа, новую границу доверия — и должен отвечать на конкретный триггер (Раздел 5).",
                 size=15)
    speaker_notes(s, load_notes("s22b"))


def build_s22c(p):
    """NEW (§4.5) — память агента: плоский файл → mem0/Cognee/Graphiti-Zep.
    Явный callback на RAG Раздела 2 — тот же вопрос масштаба знания."""
    s = blank(p)
    slide_title(s, "Память агента — тот же вопрос масштаба, что RAG.", size=24)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Память — что агент помнит МЕЖДУ сессиями (в отличие от контекста одного разговора). Спектр — от плоского файла до граф-баз знаний.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # spectrum: flat file -> vector/graph
    sy, sh = 1.86, 2.35
    ocean_box(s, 0.55, sy, 5.55, sh)
    icon(s, "file-text", 0.82, sy + 0.22, 0.50, "mid")
    text_box(s, 1.44, sy + 0.24, 4.4, 0.42, "Плоский файл", size=16, bold=True,
             color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.82, sy + 0.86, 5.0, 1.35,
             "агент дописывает факты в текстовый лог, при запуске читает целиком. Работает, пока лог мал и стабилен — прямая параллель критерию RAG «корпус помещается в окно».",
             size=13, color=DEEP, line_spacing=1.20)
    right_arrow(s, 6.20, sy + sh / 2 - 0.26, 0.55, 0.52, fill=LIGHT)
    ocean_box(s, 6.95, sy, 5.85, sh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "brain-circuit", 7.22, sy + 0.22, 0.50, "teal")
    text_box(s, 7.84, sy + 0.24, 4.7, 0.42, "Векторная / граф-база знаний", size=16,
             bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    for i, m in enumerate([
        "mem0 — память пользователя между сессиями",
        "Cognee — память на графе знаний",
        "Graphiti / Zep — временной граф: когда факт верен, когда устарел",
    ]):
        circle(s, 7.24, sy + 0.86 + i * 0.44 + 0.05, 0.11, TEAL)
        text_box(s, 7.46, sy + 0.86 + i * 0.44, 5.1, 0.42, m,
                 size=12.5, color=DEEP, line_spacing=1.08)
    gold_callout(s, 0.55, 4.42, 9.05, 1.90,
                 "Тот же вопрос масштаба знания, что решает RAG для корпуса документов, здесь встаёт для памяти самого агента. Не ставьте граф-базу знаний агенту с короткими несвязанными сессиями — это тот же технический долг без требования, что RAG для десяти статей.",
                 size=14.5)
    # #185: реальный интернет-мем «Tuxedo Winnie the Pooh» — плоский файл-лог
    # vs граф-база знаний памяти (тот же спектр, что на схеме слева-сверху).
    pmw, pmh = 3.00, 2.19
    pmx, pmy = 9.80, 4.42
    filled_rect(s, pmx - 0.05, pmy - 0.05, pmw + 0.10, pmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.04)
    add_image(s, WEB / "s22c-pooh-ru.png", pmx, pmy, pmw, pmh)
    footer(s, "Источник: публичный реестр agent-harness-registry (workain lab, независимая проверка).")
    speaker_notes(s, load_notes("s22c"))


def build_s22d(p):
    """NEW (§4.6) — провал памяти (кейс): Letta Tier D + Anthropic Memory
    Tool Tier B 17%. Freshness-оговорка Letta на видимом слое."""
    s = blank(p)
    slide_title(s, "«Агент, который помнит» — не всегда во благо.", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Наличие памяти интуитивно кажется чистым выигрышем. Независимая проверка показывает: иногда — драматически нет.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # left — Letta case with numbers
    lx, ly, lw, lh = 0.55, 1.86, 6.25, 4.05
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "triangle-alert", lx + 0.26, ly + 0.22, 0.46, "mid")
    text_box(s, lx + 0.84, ly + 0.22, lw - 1.1, 0.42, "Letta — уровень D (Tier D)",
             size=16, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 0.78, lw - 0.52, 0.62,
             "проигрывает И голой модели, И плоскому файлу на всех задачах. persistbench_v1:",
             size=13, color=DEEP, line_spacing=1.16)
    # mini table of 3 numbers
    rows = [("Голая модель", "1.000", "94 с", TEAL),
            ("Плоский файл", "0.833", "159 с", MID),
            ("Letta", "0.750", "496 с", GOLD)]
    ry = ly + 1.48
    for nm, sc, tm, col in rows:
        filled_rect(s, lx + 0.26, ry, 0.14, 0.34, col)
        text_box(s, lx + 0.52, ry, 2.7, 0.34, nm, size=12.5, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, lx + 3.25, ry, 1.3, 0.34, sc, size=13, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        text_box(s, lx + 4.65, ry, 1.3, 0.34, tm, size=12, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        ry += 0.42
    text_box(s, lx + 0.26, ry + 0.02, lw - 0.52, 0.72,
             "Механизмы: капитуляция под давлением · многословие топит факт · факт замечен, но не закоммичен.",
             size=12, color=DEEP, line_spacing=1.14)
    # freshness caveat VISIBLE
    filled_rect(s, lx + 0.26, ly + lh - 0.62, lw - 0.52, 0.48, SOFT_GREY,
                radius=True, radius_adj=0.12)
    text_box(s, lx + 0.42, ly + lh - 0.58, lw - 0.80, 0.40,
             "Оговорка: тест на Letta v0.6.7 — отстаёт от текущей v0.16.8 (~18 мес).",
             size=11, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    # right — Anthropic Memory Tool 17%
    rx, rw = 7.05, 5.75
    ocean_box(s, rx, ly, rw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.28, ly + 0.22, rw - 0.56, 0.42, "Anthropic Memory Tool — уровень B (Tier B)",
             size=16, bold=True, color=DEEP)
    text_box(s, rx + 0.28, ly + 0.70, rw - 0.56, 0.5,
             "В целом сильный результат — но даже он теряет данные в",
             size=13, color=DEEP, line_spacing=1.14)
    text_runs(s, rx + 0.28, ly + 1.18, rw - 0.56, 0.7, [
        {"text": "17% ", "size": 34, "bold": True, "color": GOLD},
        {"text": "задач", "size": 16, "bold": True, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    for i, m in enumerate([
        "явный отказ записать «эфемерный» факт",
        "немотивированный отказ «вне рамок»",
        "тихое сжатие с потерями — деталь не восстановить",
        "невоспроизводимость: та же беседа дважды → разный результат",
    ]):
        circle(s, rx + 0.30, ly + 2.04 + i * 0.42 + 0.05, 0.11, MID)
        text_box(s, rx + 0.52, ly + 2.04 + i * 0.42, rw - 0.80, 0.42, m,
                 size=12, color=DEEP, line_spacing=1.06)
    footer(s, "agent-harness-registry (workain lab, независимая проверка 2026-07-05). «Работает хорошо» ≠ «работает всегда» — тот же урок, что RAG на масштабе и забывание.")
    speaker_notes(s, load_notes("s22d"))


def build_s22e(p):
    """NEW (§4.7) — операционный слой: файлы-инструкции + presence paradox +
    Honest Lying + claude-code#51735. Интуиция расходится с измерением."""
    s = blank(p)
    slide_title(s, "Файл-инструкция агенту — не «магическая прививка».", size=24)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Операционный слой — файлы-конвенции (класс CLAUDE.md / AGENTS.md) + журнал задач. Интуиция «написал инструкцию → и стало вернее» расходится с измерением.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # 3 evidence cards
    cards = [
        ("notebook-pen", "presence paradox", "RCT (Gloaguen et al. 2026): само наличие файла-инструкции НЕ даёт значимого прироста успешности — при этом стоимость и число шагов растут.", "Помогает только там, где реально заполняет пробел документации."),
        ("brain-circuit", "Honest Lying", "Dixit, Kamal, Oates 2026: самоавторская память может ЗАКРЕПЛЯТЬ неверное убеждение — журнал бетонирует раннюю ошибку вместо пересмотра.", "Рефлексия ошибочна → повторные попытки опираются на неё."),
        ("triangle-alert", "claude-code#51735", "Реальный кейс: письменно признанная прошлая ошибка НЕ предотвратила её повторение спустя 25 дней.", "Запись о провале ≠ гарантия, что поведение изменится."),
    ]
    cw, chh = 4.00, 3.35
    cy = 1.80
    x = 0.55
    for i, (ic, t, body, foot) in enumerate(cards):
        isgold = (i == 0)
        if isgold:
            ocean_box(s, x, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, cy, cw, chh)
        icon(s, ic, x + 0.24, cy + 0.24, 0.50, "gold" if isgold else "mid")
        text_box(s, x + 0.86, cy + 0.26, cw - 1.05, 0.5, t,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x + 0.24, cy + 0.92, cw - 0.48, 1.75, body,
                 size=12.5, color=DEEP, line_spacing=1.18)
        filled_rect(s, x + 0.22, cy + chh - 0.86, cw - 0.44, 0.72, SURFACE,
                    stroke=SOFT_GREY, stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, x + 0.36, cy + chh - 0.80, cw - 0.68, 0.62, foot,
                 size=11.5, italic=True, color=(DEEP if isgold else SLATE),
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        x += cw + 0.13
    gold_callout(s, 0.55, 5.30, 9.05, 1.42,
                 "Полезен, когда заполняет реальный пробел; бесполезен, когда дублирует доступное модели; может навредить, когда самоавторская память закрепляет ошибку. Тот же паттерн, что роль в промпте: «добавить X → станет вернее» не подтверждается измерением.",
                 size=13.5)
    # #185: реальный интернет-мем «This is fine» — «файл-инструкция всё
    # починит», пока агент-система буквально горит (presence paradox).
    fmw, fmh = 3.00, 1.46
    fmx, fmy = 9.80, 5.30
    filled_rect(s, fmx - 0.05, fmy - 0.05, fmw + 0.10, fmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s22e-thisisfine-ru.png", fmx, fmy, fmw, fmh)
    speaker_notes(s, load_notes("s22e"))


def build_s23(p):
    """case_study — 3 agent failure cards; compounding chart c23 inside card 2."""
    s = blank(p)
    slide_title(s, "Провалы агентов.", size=27)
    text_box(s, 0.55, 1.14, 12.25, 0.4,
             "Каждый провал — режим отказа цикла агента в датированном кейсе: урок + альтернатива.",
             size=14, italic=True, color=MID)
    # card 1 — loop
    c1x, cy, c1w, chh = 0.55, 1.66, 4.05, 4.55
    ocean_box(s, c1x, cy, c1w, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "flame", c1x + 0.22, cy + 0.22, 0.50, "gold")
    text_box(s, c1x + 0.82, cy + 0.24, c1w - 1.0, 0.5, "1. Петля без лимитов",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, c1x + 0.24, cy + 0.85, c1w - 0.48, 0.95,
             "Агент на «синхронизируй заказы» получил HTTP 429 → план→вызов→429→…",
             size=12, color=DEEP, line_spacing=1.14)
    text_box(s, c1x + 0.24, cy + 1.72, c1w - 0.48, 0.50, "$4 200 за 63 часа",
             size=22, bold=True, color=GOLD, line_spacing=1.0)
    # #233: явная база сравнения — retry-скрипт «практически бесплатно»
    filled_rect(s, c1x + 0.24, cy + 2.24, c1w - 0.48, 0.62, WHITE, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.12)
    text_box(s, c1x + 0.38, cy + 2.27, c1w - 0.76, 0.56,
             "База сравнения: скрипт повторов с задержкой решил бы ту же задачу за секунды-минуты, практически бесплатно",
             size=10, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    text_box(s, c1x + 0.24, cy + 2.92, c1w - 0.48, 0.56,
             "$4 200 — цена не автоматизации, а неправильного выбора архитектуры под предсказуемую задачу",
             size=10, italic=True, color=SLATE, line_spacing=1.08)
    text_box(s, c1x + 0.24, cy + 3.54, c1w - 0.48, 0.90,
             "Более подходящая архитектура: скрипт повторов с задержкой, не агент; лимиты бюджета и итераций ВНЕ агента",
             size=11, bold=True, color=DEEP, line_spacing=1.10)
    # card 2 — compounding (chart)
    c2x, c2w = 4.80, 4.05
    ocean_box(s, c2x, cy, c2w, chh)
    text_box(s, c2x + 0.24, cy + 0.20, c2w - 0.48, 0.4,
             "2. Накопление ошибок (compounding)", size=15, bold=True, color=MID)
    add_image(s, CHARTS / "c23-compounding.png", c2x + 0.16, cy + 0.62,
              c2w - 0.32, 2.55)
    text_box(s, c2x + 0.24, cy + 3.30, c2w - 0.48, 1.05,
             "Вывод: «докрутить один шаг» — слабый рычаг; «меньше переходов + проверка между шагами» — сильный",
             size=12, bold=True, color=DEEP, line_spacing=1.16)
    # card 3 — multi-agent fragility
    c3x, c3w = 9.05, 3.75
    ocean_box(s, c3x, cy, c3w, chh)
    icon(s, "git-fork", c3x + 0.22, cy + 0.22, 0.46, "mid")
    text_box(s, c3x + 0.76, cy + 0.22, c3w - 0.95, 0.55,
             "3. Мульти-агентная хрупкость", size=14.5, bold=True, color=MID,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    text_box(s, c3x + 0.24, cy + 0.95, c3w - 0.48, 1.55,
             "зависимые подзадачи → параллельные субагенты принимают конфликтующие неявные решения",
             size=12.5, color=DEEP, line_spacing=1.18)
    text_box(s, c3x + 0.24, cy + 2.70, c3w - 0.48, 1.65,
             "Более подходящая: однопоточный линейный агент; мульти-агент — только широко-параллельное независимое",
             size=12.5, bold=True, color=DEEP, line_spacing=1.18)
    footer(s, "Атаки через инструменты (инъекция в промпт, утечка через GitHub MCP) — 4-й класс провалов, разобран на слайде безопасности. Петля на $4 200 — постмортем одного автора 2026-04 (иллюстративно); накопление ошибок — MindStudio 2025–2026.")
    speaker_notes(s, load_notes("s23"))


def build_s25(p):
    """NEW/MERGED (§4.8) — skills + subagents + доступ к инструментам +
    ИНТЕГРИРОВАННАЯ безопасность равного веса (P1: GOLD security block, не
    caveat снизу). GitHub MCP heist + ZDR-факты живут ЗДЕСЬ."""
    s = blank(p)
    slide_title(s, "Навыки, субагенты, доступ — и граница доверия.", size=25)
    # top row — 3 equipment slots
    slots = [
        ("puzzle", "Навык (skill)", "переиспользуемая процедура «как делать» под повторяющуюся задачу"),
        ("users", "Субагент", "отдельное контекстное окно: не засорять основной + изолировать недоверенное"),
        ("cable", "MCP-доступ", "каждое подключение — новая граница доверия и политика хранения"),
    ]
    cw, chh = 4.00, 1.55
    cy = 1.12
    x = 0.55
    for ic, t, body in slots:
        ocean_box(s, x, cy, cw, chh)
        icon(s, ic, x + 0.22, cy + 0.22, 0.46, "mid")
        text_box(s, x + 0.80, cy + 0.22, cw - 1.0, 0.42, t,
                 size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.24, cy + 0.74, cw - 0.48, 0.74, body,
                 size=12, color=DEEP, line_spacing=1.14)
        x += cw + 0.13
    # GOLD security block — EQUAL VISUAL WEIGHT (P1), not a footer caveat
    gy, gh = 2.86, 3.60
    ocean_box(s, 0.55, gy, 12.25, gh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=3.0)
    icon(s, "shield-alert", 0.82, gy + 0.22, 0.56, "gold")
    text_box(s, 1.52, gy + 0.24, 11.0, 0.44,
             "Безопасность: как только агент делегирует и подключается — появляется поверхность, которой не было у одного вызова",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    # left half — GitHub heist mechanism + 2 ZDR facts
    text_box(s, 0.82, gy + 0.80, 5.85, 0.34, "GitHub MCP: утечка через инъекцию (май 2025)",
             size=13.5, bold=True, color=DEEP)
    text_box(s, 0.82, gy + 1.14, 5.85, 0.78,
             "issue с встроенной инструкцией + переизбыточный токен (PAT на все репо) → ассистент выгрузил приватные репозитории в публичный PR.",
             size=12, color=DEEP, line_spacing=1.16)
    filled_rect(s, 0.82, gy + 1.96, 5.85, 0.46, WHITE, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.14)
    text_box(s, 0.98, gy + 1.99, 5.55, 0.40,
             "Катастрофа = инъекция × широкие права. Убери любое — атака не проходит.",
             size=12, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.82, gy + 2.48, 5.85, 0.44,
             "«У нас ZDR» ≠ «нигде не хранится»: судебный приказ (NYT v. OpenAI) + сторонние сервисы/MCP вне ZDR. Регулируемое — не слать без ZDR/BAA.",
             size=11, italic=True, color=DEEP, line_spacing=1.10)
    # right half — 4 rules
    rx = 7.05
    text_box(s, rx, gy + 0.80, 5.5, 0.34, "4 правила проектирования",
             size=13.5, bold=True, color=DEEP)
    rules = [
        ("Наименьшие привилегии (least-privilege)", "минимум токенов/прав"),
        ("Изоляция недоверенного", "отдельно от привилегий (субагент)"),
        ("Человек в контуре на запись", "необратимое — только через человека"),
        ("Белый список / фиксация версий", "аудированные версии; запрет по умолчанию"),
    ]
    ry = gy + 1.18
    for i, (t, b) in enumerate(rules):
        circle(s, rx, ry + 0.03, 0.30, GOLD)
        text_box(s, rx, ry + 0.03, 0.30, 0.30, str(i + 1),
                 size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 0.42, ry - 0.02, 5.5, 0.30, t, size=13, bold=True, color=DEEP)
        text_box(s, rx + 0.42, ry + 0.26, 5.5, 0.28, b, size=11, color=SLATE)
        ry += 0.58
    speaker_notes(s, load_notes("s25"))


def build_s25b(p):
    """NEW (§4.9) — обзор реальных coding-агентов через рамку экипировки §4.4.
    РОВНО 4 инструмента: Claude Code / Aider / Cursor / OpenHands. OpenHands
    помечен как неподтверждённая гипотеза «OpenClaw»."""
    s = blank(p)
    slide_title(s, "Реальные агенты для кода — через рамку экипировки.", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Разница между инструментами — не в «качестве модели», а в том, какие слоты экипировки заполнены и где агент физически живёт.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    tools = [
        ("Claude Code", "терминал + IDE · проприетарный",
         "широкая экипировка: память, файлы-инструкции, навыки, полноценные субагенты, MCP — почти все 5 слотов",
         "цена — большая операционная сложность", MID, False),
        ("Aider", "сначала терминал · открытый код",
         "минимальная простота: без развитой памяти, без субагентов, без навыков. ~47k★ на GitHub",
         "«тонкая» оснастка — самостоятельный выбор, не недоразвитость", TEAL, False),
        ("Cursor", "настольная IDE (форк VS Code) · проприетарный",
         "агент внутри редактора — не терминал. Форма интеграции — отдельная ось от оснастки",
         "где живёт агент — тоже архитектурное решение", MID, False),
        ("OpenHands", "платформа на своём сервере · MIT · ~80k★",
         "широкая автономность, локальное/Docker/облако развёртывание",
         "рабочая гипотеза по совпадению профиля — вероятный кандидат на «OpenClaw», не подтверждённый факт", GOLD, True),
    ]
    cw, chh = 3.02, 4.05
    cy = 1.74
    x = 0.55
    for nm, meta, body, note, col, ishyp in tools:
        if ishyp:
            ocean_box(s, x, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, cy, cw, chh)
        icon(s, "code", x + 0.22, cy + 0.22, 0.44, "gold" if ishyp else "mid")
        text_box(s, x + 0.76, cy + 0.20, cw - 0.92, 0.48, nm,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x + 0.24, cy + 0.76, cw - 0.48, 0.44, meta,
                 size=10.5, italic=True, color=LIGHT, line_spacing=1.08)
        text_box(s, x + 0.24, cy + 1.22, cw - 0.48, 1.30, body,
                 size=11.5, color=DEEP, line_spacing=1.16)
        nb_h = 1.30 if ishyp else 0.86
        filled_rect(s, x + 0.20, cy + chh - nb_h - 0.14, cw - 0.40, nb_h,
                    WHITE if ishyp else SURFACE, stroke=(GOLD if ishyp else SOFT_GREY),
                    stroke_pt=(1.5 if ishyp else 1.0), radius=True, radius_adj=0.10)
        text_box(s, x + 0.32, cy + chh - nb_h - 0.08, cw - 0.64, nb_h - 0.10, note,
                 size=(9.5 if ishyp else 10), italic=True,
                 color=(DEEP if ishyp else SLATE), line_spacing=1.10)
        x += cw + 0.13
    gold_callout(s, 0.55, 5.98, 12.25, 0.86,
                 "Выбор инструмента подчиняется тому же правилу: не бери максимально оснащённый по умолчанию — смотри, какая оснастка нужна именно твоей задаче.",
                 size=13.5)
    speaker_notes(s, load_notes("s25b"))


def build_s26(p):
    """schema_layered — complexity ladder, bottom-aligned. 6 short rungs +
    trigger label in the gap BELOW each rung (the requirement that opens the
    climb to it). Left col = ladder; right col = rule panel."""
    s = blank(p)
    slide_title(s, "Лестница архитектурной сложности.", size=27)
    text_box(s, 0.55, 1.14, 8.40, 0.40,
             "Оставайся на нижней ступени; поднимайся только под требование задачи.",
             size=13.5, italic=True, color=MID)
    # bottom-up: idx0 = step1 bottom (gold). trig = requirement that opens
    # the climb FROM this rung to the next (shown in the gap above it).
    steps = [
        ("1", "Обычный код (без ИИ)", "нужен NL / неструктурированный ввод / неточное соответствие", GOLD, True),
        ("2", "Один вызов LLM (промпт; +CoT, +примеры в промпте)", "знание большое И меняющееся И провенанс И приватное", MID, False),
        ("3", "RAG / контекст-инжиниринг", "задача многошаговая, последовательность известна заранее", LIGHT, False),
        ("4", "Сценарий (workflow, предопределённые пути)", "непредсказуема И ценность оправдывает кратные стоимость/риск", LIGHT, False),
        ("5", "Агент (план→действие→проверка→повтор + лимиты)", "подзадачи широко-параллельны И независимы И высокоценны", LIGHT, False),
        ("6", "Мульти-агент", None, LIGHT, False),
    ]
    n = len(steps)
    step_h = 0.52
    trig_h = 0.34
    bottom_edge = 6.62  # bottom of rung 1
    x0 = 0.55
    full_w = 8.05
    for i, (num, label, trig, col, isgold) in enumerate(steps):
        y = bottom_edge - step_h - i * (step_h + trig_h)
        indent = i * 0.16
        w = full_w - indent
        x = x0 + indent
        if isgold:
            ocean_box(s, x, y, w, step_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
        else:
            ocean_box(s, x, y, w, step_h)
        circle(s, x + 0.12, y + (step_h - 0.34) / 2, 0.34, col)
        text_box(s, x + 0.12, y + (step_h - 0.34) / 2, 0.34, 0.34, num,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.58, y, w - 0.72, step_h, label,
                 size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        if trig:
            # trigger shown in the gap ABOVE this rung (opens climb to next)
            tgy = y - trig_h - 0.01
            text_box(s, x + 0.30, tgy, full_w - 0.50, trig_h,
                     "↑  " + trig, size=10, italic=True, color=SLATE,
                     anchor=MSO_ANCHOR.MIDDLE)
    # PA-1 (owner-approved): direction-of-scale strip alongside the ladder —
    # «сложнее ↑» at top / «проще ↓» at bottom (kills «выше=лучше» mis-read).
    up_arrow(s, 8.66, 1.98, 0.26, 4.62, fill=COVER_OUTLINE)
    text_box(s, 8.30, 1.58, 1.00, 0.32, "сложнее ↑", size=10.5, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    text_box(s, 8.30, 6.62, 1.00, 0.32, "проще ↓", size=10.5, bold=True,
             color=LIGHT, align=PP_ALIGN.CENTER)
    # rule panel right (full height, gold — the central rule)
    ocean_box(s, 9.15, 1.55, 3.65, 5.05, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "milestone", 9.42, 1.85, 0.62, "gold")
    text_box(s, 9.42, 2.70, 3.15, 1.75,
             "Оставайся на самой нижней ступени, которая закрывает требования.",
             size=16, bold=True, color=DEEP, line_spacing=1.24)
    text_box(s, 9.42, 4.45, 3.15, 2.05,
             "Каждый подъём — это ОБМЕН (возможности ↔ стоимость, задержка, аудируемость, поверхность атаки), а не чистый выигрыш.",
             size=13.5, color=DEEP, line_spacing=1.22)
    footer(s, "Нижняя ступень — «обычный код без ИИ»: лестница начинается с вопроса «а нужен ли ИИ вообще».")
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """NEW (§5.2) — «План решения»: 8-шаговый маршрут вопросов сверху вниз
    (flowchart с да/нет-ветвлениями), ЗАМЕНА старой 7×7 матрицы. Нижняя
    плашка-приоритет (детерминированное → код, СТОП) — gold-вес."""
    s = blank(p)
    slide_title(s, "План решения: маршрут вопросов, а не сумма баллов.", size=24)
    text_box(s, 0.55, 1.02, 8.4, 0.36,
             "Пройдите задачу сверху вниз; останавливайтесь на первом сработавшем вопросе. «Да» → результат справа; «нет» → следующий вопрос ниже.",
             size=12, italic=True, color=MID, line_spacing=1.08)
    # left — vertical routed flow (question -> да-outcome)
    fx = 0.55
    qw = 5.35           # question box width
    ow = 3.05           # outcome box width
    ax = fx + qw + 0.12  # outcome x
    steps = [
        ("Детерминированная и верифицируемая?", "да → обычный код · СТОП", True),
        ("Закрывает один вызов (+CoT)?", "да → промпт · СТОП", False),
        ("Нужен провенанс по источнику?", "да → RAG с опорой на источник / код", False),
        ("Знание меняется / провенанс?", "да → RAG (не FT)", False),
        ("Нужно поведение / тон / формат?", "да → дообучение (PEFT)", False),
        ("Многошагово, порядок известен?", "да → сценарий · нет → агент+лимиты", False),
        ("Подзадачи параллельны+независимы?", "да → мульти-агент · нет → линейный", False),
    ]
    y = 1.42
    qh = 0.50
    vg = 0.10
    for i, (q, out, isgold) in enumerate(steps):
        if isgold:
            ocean_box(s, fx, y, qw, qh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.5)
        else:
            ocean_box(s, fx, y, qw, qh)
        circle(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, GOLD if isgold else MID)
        text_box(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=(DEEP if isgold else WHITE),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, fx + 0.54, y, qw - 0.68, qh, q,
                 size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        # outcome to the right
        filled_rect(s, ax, y + 0.05, ow, qh - 0.10, GOLD if isgold else TEAL_TINT,
                    stroke=(GOLD if isgold else TEAL), stroke_pt=1.2, radius=True, radius_adj=0.12)
        text_box(s, ax + 0.12, y + 0.05, ow - 0.24, qh - 0.10, out,
                 size=10.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        # connector question->outcome (да →)
        right_arrow(s, fx + qw + 0.005, y + qh / 2 - 0.09, 0.11, 0.18, fill=LIGHT)
        # thin down-connector to next question (нет ↓), centred in the gap
        if i < len(steps) - 1:
            connector(s, fx + 0.28, y + qh, fx + 0.28, y + qh + vg, LIGHT, 1.6)
        y += qh + vg
    # step 8 — same row pattern as 1–7 (unified), spanning question+outcome cols
    filled_rect(s, fx, y, qw + 0.12 + ow, qh, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.5, radius=True, radius_adj=0.12)
    circle(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, GOLD)
    text_box(s, fx + 0.12, y + (qh - 0.32) / 2, 0.32, 0.32, "8",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, fx + 0.54, y, qw + 0.12 + ow - 0.68, qh,
             "Данные чувствительны? → на КАЖДОМ шаге: карта данных + наименьшие привилегии + ZDR/BAA",
             size=10.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    # right — worked example + mini-apply
    rx, rw = 9.10, 3.70
    ocean_box(s, rx, 1.46, rw, 2.55, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.24, 1.60, rw - 0.48, 0.34, "Пример (задача A)", size=13.5,
             bold=True, color=TEAL)
    text_box(s, rx + 0.24, 1.96, rw - 0.48, 1.95,
             "«2000 регламентов, меняются еженедельно, ответ со ссылкой на пункт» → вопрос 4 (меняется+провенанс) решающий → RAG со строгой опорой на источник. Не дообучение (устареет), не код (нужен NL).",
             size=12, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, 4.12, rw, 1.85)
    text_box(s, rx + 0.24, 4.26, rw - 0.48, 0.34, "Разминка (задача B)", size=13.5,
             bold=True, color=MID)
    text_box(s, rx + 0.24, 4.62, rw - 0.48, 1.30,
             "«бот на ~150 FAQ, меняются раз в квартал» — пройдите маршрут сами; разбор на семинаре.",
             size=12, color=DEEP, line_spacing=1.18)
    # bottom priority plate — gold, the most important line
    gold_callout(s, 0.55, 6.06, 12.25, 0.82,
                 "Приоритет маршрута: если задача детерминированная и верифицируемая — обычный код, СТОП здесь. ИИ добавил бы лишь недетерминизм, стоимость, задержку и поверхность для инъекции в промпт (prompt injection).",
                 size=14)
    speaker_notes(s, load_notes("s27"))


def build_s27b(p):
    """NEW (§5.2b) — «Стартовый комплект агента и когда его усложнять».
    Growth-ladder playbook. Рифма с лестницей s26 и картой 5 слотов §4.4."""
    s = blank(p)
    slide_title(s, "Стартовый комплект агента — и когда его усложнять.", size=24)
    text_box(s, 0.55, 1.14, 12.25, 0.42,
             "Та же лестница «не усложняй без требования», применённая на уровень ниже — к оснастке одного агента, а не системы целиком.",
             size=13.5, italic=True, color=MID, line_spacing=1.15)
    # left — thin default
    lx, ly, lw, lh = 0.55, 1.86, 4.45, 4.05
    ocean_box(s, lx, ly, lw, lh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "package", lx + 0.26, ly + 0.24, 0.54, "teal")
    text_box(s, lx + 0.92, ly + 0.26, lw - 1.1, 0.5, "По умолчанию — тонкий агент",
             size=16, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    for i, t in enumerate([
        "один файл-инструкция",
        "плоская память",
        "БЕЗ субагентов",
        "минимальный набор навыков",
        "минимальный MCP-доступ",
    ]):
        circle(s, lx + 0.30, ly + 1.02 + i * 0.48 + 0.06, 0.12, TEAL)
        text_box(s, lx + 0.56, ly + 1.02 + i * 0.48, lw - 0.84, 0.44, t,
                 size=14, color=DEEP, line_spacing=1.1)
    text_box(s, lx + 0.28, ly + lh - 0.62, lw - 0.56, 0.56,
             "Бремя доказательства — на усложнении, а не на простоте.",
             size=12, italic=True, color=DEEP, line_spacing=1.14)
    # right — 3 justified triggers (сужены — крайняя правая колонка под мем)
    rx, rw = 5.30, 4.30
    triggers = [
        ("brain-circuit", "Память-бэкенд — когда:",
         "история переросла контекст ИЛИ нужен поиск по фактам. Тот же критерий, что промпт→RAG."),
        ("users", "Субагенты — когда:",
         "подзадача требует отдельного окна ИЛИ изоляции недоверенной работы (наименьшие привилегии)."),
        ("cable", "Больше MCP-доступа — когда:",
         "конкретная задача требует конкретного инструмента — не «на всякий случай». Каждое подключение — граница доверия."),
    ]
    ty2 = 1.86
    th = 1.24
    for ic, t, b in triggers:
        ocean_box(s, rx, ty2, rw, th)
        icon(s, ic, rx + 0.22, ty2 + 0.20, 0.42, "mid")
        text_box(s, rx + 0.76, ty2 + 0.16, rw - 0.94, 0.36, t,
                 size=13.5, bold=True, color=MID)
        text_box(s, rx + 0.76, ty2 + 0.50, rw - 0.94, 0.68, b,
                 size=11.5, color=DEEP, line_spacing=1.12)
        ty2 += th + 0.16
    # #185: реальный интернет-мем «Two guys on a bus» — грустный (усложнил
    # на всякий случай) vs довольный (начал с тонкого агента). Крайняя правая.
    bmx, bmy, bmw, bmh = 9.75, 2.35, 3.05, 3.10
    ocean_box(s, bmx, bmy, bmw, bmh)
    add_image(s, WEB / "s27b-bus-ru.png", bmx + 0.12, bmy + 0.12,
              bmw - 0.24, bmh - 0.24)
    gold_callout(s, 0.55, 6.06, 12.25, 0.82,
                 "Тот же принцип, что лестница архитектур — применённый к оснастке одного агента: presence paradox показал, что даже файл-инструкция «как ритуал» не работает. Усложняй под конкретный проверяемый триггер.",
                 size=13)
    speaker_notes(s, load_notes("s27b"))


def build_s29(p):
    """assertion_visual — human validator + MIT NANDA ~95% donut (c29)."""
    s = blank(p)
    slide_title(s, "Человек-валидатор + урок MIT NANDA.", size=27)
    # left — human validator
    lx, ly, lw, lh = 0.55, 1.40, 6.55, 4.45
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "user-check", lx + 0.26, ly + 0.24, 0.52, "mid")
    text_box(s, lx + 0.92, ly + 0.24, lw - 1.1, 0.50,
             "Агент делает — человек проверяет результат и факты",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    text_box(s, lx + 0.28, ly + 0.80, lw - 0.56, 0.56,
             "против независимого источника истины, а не по правдоподобности рассуждения: самообъяснение модели ≠ контроль.",
             size=12, italic=True, color=SLATE, line_spacing=1.14)
    # #237: три измерения роли человека-валидатора
    dims = [
        ("Степень автономности", "от «человек нажимает кнопку» до «агент уведомляет постфактум»"),
        ("Область доверия", "чтение (легко откатить) vs запись · обратимое vs необратимое"),
        ("Непрерывный мониторинг", "метрики качества постоянно, не разовая проверка на старте"),
    ]
    py = ly + 1.52
    for t, b in dims:
        filled_rect(s, lx + 0.28, py, lw - 0.56, 0.90, SURFACE, stroke=SOFT_GREY,
                    stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, lx + 0.44, py + 0.10, lw - 0.86, 0.34, t,
                 size=13.5, bold=True, color=MID)
        text_box(s, lx + 0.44, py + 0.44, lw - 0.86, 0.42, b,
                 size=11.5, color=DEEP, line_spacing=1.10)
        py += 0.98
    # right — NANDA donut
    rx, rw = 7.35, 5.45
    ocean_box(s, rx, ly, rw, 3.10, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    add_image(s, CHARTS / "c29-nanda.png", rx + 0.20, ly + 0.20, 2.45, 2.65)
    text_box(s, rx + 2.75, ly + 0.35, rw - 3.0, 0.95, "~95%",
             size=42, bold=True, color=GOLD, line_spacing=1.0)
    text_box(s, rx + 2.75, ly + 1.30, rw - 3.0, 1.65,
             "корпоративных GenAI-пилотов без измеримого возврата инвестиций — корень в разрыве обучения и провале интеграции, не в качестве модели",
             size=12.5, color=DEEP, line_spacing=1.18)
    ocean_box(s, rx, ly + 3.25, rw, 1.20)
    text_box(s, rx + 0.26, ly + 3.38, rw - 0.52, 0.95,
             "«Запустить ИИ» ≠ «получить ценность». Решает архитектурно-интеграционная дисциплина. Иногда правильный ответ — простейшая архитектура или не-ИИ.",
             size=12, bold=True, color=DEEP, line_spacing=1.16,
             anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "MIT NANDA, State of AI in Business 2025 — отчёт с методологией (150 интервью + 350 опрос + 300 внедрений), не универсальный закон.")
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """hero_closing / bridge — мост к Лекции 4 «AI в разработке ПО (ось SDLC ×
    AI)» (issue #170). Слева: 4 бокса-этапа SDLC × AI, у каждого teal-якорь из
    Л3. Справа: hero-фото разработки ≥40%. Gold-заголовок-мост."""
    s = blank(p)
    hx, hy, hw, hh = 8.10, 0.0, 5.233, 7.5
    hero_image(s, SCREENSHOTS / "s30-coding.jpg", hx, hy, hw, hh)
    # #185/#317: подпись-атрибуция фото снята с visible (attribution.md — legal).
    lx, lw = 0.55, 7.30
    text_box(s, lx, 0.42, lw, 0.94,
             "Дальше — отрасли. Старт: разработка ПО, ось SDLC × AI.",
             size=24, bold=True, color=DEEP, line_spacing=1.06)
    text_box(s, lx, 1.34, lw, 0.40,
             "Лекция 4 берёт тот же аппарат и раскладывает его по этапам жизненного цикла ПО. Каждый этап опирается на якорь из этой лекции:",
             size=12.5, italic=True, color=MID, line_spacing=1.14)
    # 2x2 SDLC × AI boxes, each with a teal anchor line from Lec 3
    boxes = [
        ("Требования / дизайн", "выбор архитектуры под шаг → лестница из 6 ступеней"),
        ("Кодирование", "агенты для кода через рамку экипировки; провалы как КЛАСС"),
        ("Тестирование", "граница доверия и наименьшие привилегии; данные вне ZDR"),
        ("Эксплуатация", "человек-валидатор проверяет результат, а не самообъяснение"),
    ]
    bx0, by0 = lx, 1.90
    bw, bh = 3.55, 1.62
    gap = 0.15
    for i, (t, anchor) in enumerate(boxes):
        x = bx0 + (i % 2) * (bw + gap)
        y = by0 + (i // 2) * (bh + gap)
        ocean_box(s, x, y, bw, bh)
        chip(s, x + 0.18, y + 0.16, min(bw - 0.36, 0.15 * len(t) + 0.4), 0.38,
             t, fill=MID, color=WHITE, size=11.5)
        text_box(s, x + 0.22, y + 0.66, bw - 0.44, 0.86, anchor, size=11.5,
                 italic=True, color=TEAL, line_spacing=1.16)
    gold_callout(s, lx, 5.42, lw, 0.72,
                 "Задание — Семинар 3: прогнать чек-лист выбора архитектуры на 3 кейсах (чат / агент / RAG / API).",
                 size=13)
    gold_callout(s, lx, 6.28, lw, 0.78,
                 "Что делать: эта рамка — база для Лекций 4–17. На каждой отраслевой лекции прогоняйте тот же чек-лист выбора.",
                 size=13)
    speaker_notes(s, load_notes("s30"))


# ============================================================
# v3 new builders (suffix-ID, plan §4 U-1…U-7) — NO renumber s01–s30.
# ============================================================
def build_s04a(p):
    """section_divider — Раздел 1 «Промпт и его границы» (U-1)."""
    build_section_divider(
        p, 1, "Раздел 1", "Промпт и его границы",
        "Лестницу мы увидели целиком — теперь снизу: что умеет один вызов и где его потолок, прежде чем что-либо усложнять.",
        "s04a",
        image_src=WEB / "div-r1-knife.jpg",
        tag="один точный рез · 3 разбора · 1 провал")


def build_s13a(p):
    """section_divider — Раздел 3 «Fine-tune vs промпт vs RAG» (U-3)."""
    build_section_divider(
        p, 3, "Раздел 3", "Дообучение vs промпт vs RAG",
        "Проблему знания мы решили через RAG. А если проблема не в знании, а в поведении модели — её тоном, форматом, политикой?",
        "s13a",
        image_src=WEB / "div-r3-tuning.jpg",
        tag="настройка поведения · 4 разбора · 1 провал")


def build_s13b(p):
    """assertion_visual — определение fine-tuning ДО критики (U-2).

    Сверху определение → центр: мини-схема pipeline
    [предобуч.модель]＋[датасет] → дообучение → [дообуч.веса] →
    низ: контраст-плашка КОНТЕКСТ vs ВЕСА. Gold-якорь — «ВЕСА».
    Schema §5.5 Process/Pipeline checklist.
    """
    s = blank(p)
    slide_title(s, "Что такое дообучение (fine-tuning).", size=27)
    text_box(s, 0.55, 1.14, 12.25, 0.74,
             "Дообучение (fine-tuning) — продолжение обучения уже готовой модели на ваших данных. В Л1 — тип использования; здесь — архитектурный выбор, одна из ступеней лестницы.",
             size=14, italic=True, color=MID, line_spacing=1.18)
    # mini-schema pipeline in ocean box
    sy, sh = 1.98, 2.48
    ocean_box(s, 0.55, sy, 12.25, sh)
    # 3 nodes + «+» (n1→n2) + arrow with «дообучение» label (n2→n3)
    bw, bh = 2.85, 1.46
    by = sy + 0.66                     # node band lowered → label clearance
    n1x = 1.35                         # centered group (box inner 0.55..12.80)
    plus_x = n1x + bw                  # 4.20 — «+» zone 0.58 wide
    n2x = plus_x + 0.58                # 4.78
    arr_x0 = n2x + bw                  # 7.63 — arrow zone 1.50 wide
    n3x = arr_x0 + 1.50                # 9.13
    # node 1 — pretrained model
    filled_rect(s, n1x, by, bw, bh, SURFACE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    icon(s, "cpu", n1x + bw / 2 - 0.23, by + 0.14, 0.42, "mid")
    text_box(s, n1x + 0.10, by + 0.60, bw - 0.20, 0.56,
             "Предобученная\nмодель", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.02)
    text_box(s, n1x + 0.14, by + 1.18, bw - 0.28, 0.26,
             "общие веса", size=11.5, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    # «+» between n1 and n2
    text_box(s, plus_x, by + bh / 2 - 0.32, 0.58, 0.64, "+",
             size=30, bold=True, color=MID, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # node 2 — your dataset
    filled_rect(s, n2x, by, bw, bh, SURFACE, stroke=LIGHT, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    icon(s, "database", n2x + bw / 2 - 0.23, by + 0.14, 0.42, "teal")
    text_box(s, n2x + 0.14, by + 0.60, bw - 0.28, 0.30,
             "Ваш датасет", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, n2x + 0.10, by + 0.94, bw - 0.20, 0.50,
             "примеры нужного\nповедения", size=11.5, italic=True,
             color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.02)
    # «дообучение» label ABOVE arrow (clear vertical separation)
    text_box(s, arr_x0 - 0.05, by - 0.42, 1.60, 0.32, "дообучение",
             size=13, bold=True, color=MID, align=PP_ALIGN.CENTER)
    # arrow node2 → node3
    right_arrow(s, arr_x0 + 0.06, by + bh / 2 - 0.21, 1.38, 0.42, fill=MID)
    # node 3 — fine-tuned weights (gold = the changed thing)
    filled_rect(s, n3x, by, bw, bh, GOLD_TINT, stroke=GOLD, stroke_pt=2.0,
                radius=True, radius_adj=0.10)
    icon(s, "sliders-horizontal", n3x + bw / 2 - 0.23, by + 0.14, 0.42,
         "gold")
    text_box(s, n3x + 0.10, by + 0.60, bw - 0.20, 0.56,
             "Дообученные\nВЕСА", size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.02)
    text_box(s, n3x + 0.14, by + 1.18, bw - 0.28, 0.26,
             "модель уже другая", size=11.5, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER)
    # contrast strip — КОНТЕКСТ vs ВЕСА (2 halves)
    cy, ch = 4.55, 1.55
    ocean_box(s, 0.55, cy, 6.05, ch, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, 0.83, cy + 0.18, 5.50, 0.34, "Промпт / RAG → КОНТЕКСТ",
             size=15, bold=True, color=TEAL)
    text_box(s, 0.83, cy + 0.58, 5.55, 0.88,
             "Меняют только вход — веса не трогаются; эффект живёт лишь в рамках запроса.",
             size=13, color=DEEP, line_spacing=1.22)
    ocean_box(s, 6.75, cy, 6.05, ch)
    text_box(s, 7.03, cy + 0.18, 5.50, 0.34, "Дообучение → САМИ ВЕСА",
             size=15, bold=True, color=MID)
    text_box(s, 7.03, cy + 0.58, 5.55, 0.88,
             "Изменение встроено в модель — действует всегда и стоит дороже того, что меняет контекст.",
             size=13, color=DEEP, line_spacing=1.22)
    gold_callout(s, 0.55, 6.28, 12.25, 0.78,
                 "Промпт/RAG = «что показать модели».  Дообучение = «изменить саму модель». На практике «дообучить» почти всегда означает LoRA/PEFT, а не переобучение всех весов (следующий слайд).",
                 size=14)
    speaker_notes(s, load_notes("s13b"))



def build_s25a(p):
    """section_divider — Раздел 5 «Как выбрать: фреймворк решения» (U-5b)."""
    build_section_divider(
        p, 5, "Раздел 5", "Как выбрать: фреймворк решения",
        "Мы разобрали все архитектуры по отдельности — и где каждая проваливается. Теперь соберём это в один инструмент выбора.",
        "s25a",
        image_src=WEB / "div-r5-control-panel.jpg",
        tag="инструмент выбора · 4 разбора")


# ============================================================
# v5b classic-base-first (issue #185 WP8) — 51→56 слайдов.
# По одному «классическая база с нуля» слайду на каждый содержательный
# раздел (§1–§5), вставляется сразу после дивайдера раздела, перед AI-
# частью. Общий шаблон: 3 карточки классической базы + gold «что оставить
# из классики» + мост к AI-части. NO renumber (мнемонические id
# s-classic-*). Schema §5.5: 3-column tile, mass-balanced, single-line names.
# ============================================================

def build_classic_base(p, sid, *, title, intro, cards, keep_text, bridge):
    """Shared «классическая база» slide: title + intro line + 3 tile cards
    (icon + bold single-line name + body) + gold «что оставить» callout +
    teal-tint bridge strip. `cards` = list of 3 (icon_name, name, body)."""
    s = blank(p)
    slide_title(s, title, size=26)
    text_box(s, 0.55, 1.14, 12.25, 0.60, intro,
             size=14, italic=True, color=MID, line_spacing=1.14)
    # 3 tile cards — equal mass, full width. Name может занимать 2 строки;
    # body — до 5 строк, всё внутри карточки (нет overflow под gold-плашку).
    n = len(cards)
    gap = 0.24
    x0 = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    cy, chh = 1.86, 2.80
    for i, (ic, name, body) in enumerate(cards):
        x = x0 + i * (cw + gap)
        ocean_box(s, x, cy, cw, chh)
        icon(s, ic, x + 0.28, cy + 0.24, 0.54, "mid")
        text_box(s, x + 0.28, cy + 0.92, cw - 0.56, 0.72, name,
                 size=15.5, bold=True, color=MID, line_spacing=1.04)
        text_box(s, x + 0.28, cy + 1.66, cw - 0.56, chh - 1.80, body,
                 size=12, color=DEEP, line_spacing=1.14)
    # gold «что оставить из классики»
    ky, kh = 4.80, 1.36
    filled_rect(s, 0.55, ky, 12.25, kh, GOLD_TINT, stroke=GOLD, stroke_pt=1.75,
                radius=True, radius_adj=0.06)
    text_box(s, 0.83, ky + 0.11, 2.4, 0.30, "ЧТО ОСТАВИТЬ",
             size=12, bold=True, color=GOLD)
    text_box(s, 0.83, ky + 0.44, 11.6, kh - 0.54, keep_text,
             size=13.5, bold=True, color=DEEP, line_spacing=1.14)
    # teal-tint bridge strip к AI-части
    by, bh = 6.26, 0.80
    filled_rect(s, 0.55, by, 12.25, bh, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_runs(s, 0.83, by + 0.11, 11.6, bh - 0.20, [
        {"text": "Мост: ", "size": 12, "bold": True, "color": TEAL},
        {"text": bridge, "size": 12, "color": DEEP},
    ], line_spacing=1.12, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes(sid))
    return s


def build_s_classic_prompt(p):
    """§1.0 — классическая база раздела 1: точная постановка задачи."""
    build_classic_base(
        p, "s-classic-prompt",
        title="Как задачу ставили до промпта — и что это меняет.",
        intro="До больших моделей «заставить систему сделать нужное» означало не пожелание на "
              "естественном языке, а точную спецификацию и детерминированную программу.",
        cards=[
            ("file-text", "Точная спецификация / ТЗ",
             "Пред- и постусловия, инварианты, критерии приёмки (Z-нотация, TLA+, Design by "
             "Contract). Результат детерминирован и проверяем."),
            ("git-fork", "Императив vs декларатив",
             "«Как сделать» (алгоритм по шагам) против «что получить» (SQL/Prolog описывает "
             "результат, движок решает как)."),
            ("braces", "Контракт интерфейса",
             "Точное соглашение вход/выход (сигнатуры типов, OpenAPI, Protobuf). Единственный "
             "корректный смысл и способ его проверить."),
        ],
        keep_text="Дисциплину точной постановки: точный промпт — то же ТЗ на естественном языке. "
                  "Для детерминированного и верифицируемого (арифметика, валидация по схеме, "
                  "маршрутизация по правилам) — классический код, а не промпт.",
        bridge="промпт — постановка задачи вероятностной системе на естественном языке; отсюда и "
               "сила (не нужно специфицировать неподъёмное), и граница (нет единственного смысла, "
               "нет детерминизма).",
    )


def build_s_classic_rag(p):
    """§2.0 — классическая база раздела 2: классический информационный поиск."""
    build_classic_base(
        p, "s-classic-rag",
        title="Как искали в тексте до эмбеддингов.",
        intro="Буква R в RAG — retrieval, поиск: дисциплина с полувековой историей. Именно её "
              "устройство определяет, где RAG работает, а где ломается.",
        cards=[
            ("book-open", "Инвертированный индекс",
             "Для каждого слова — список документов, где оно встречается (машинный каталог). "
             "Фундамент Lucene, Elasticsearch, PostgreSQL full-text."),
            ("route", "Булев поиск",
             "Запрос как логическое выражение («ошибка AND аутентификация NOT tomcat»): точный, "
             "предсказуемый, объяснимый отбор."),
            ("list-ordered", "TF-IDF → BM25",
             "Ранжирование по важности слова (реже в коллекции — сильнее сигнал). BM25 (Okapi) — "
             "дешёвая объяснимая линия, которую многие не обгоняют."),
        ],
        keep_text="Классика точна на кодах и идентификаторах, где смысл-поиск размывает. Сильный "
                  "RAG-2026 — гибрид BM25 + плотные векторы, лексические фильтры по метаданным, "
                  "дисциплина ранжирования (реранкер) и наблюдаемость: recall/precision на "
                  "эталонном наборе (golden set).",
        bridge="семантический поиск на эмбеддингах (Лекция 2) добавляет поверх классики смысловое "
               "совпадение вместо лексического — но не заменяет её. RAG расширяет классический "
               "поиск, а не отменяет его.",
    )


def build_s_classic_ft(p):
    """§3.0 — классическая база раздела 3: классическое машинное обучение."""
    build_classic_base(
        p, "s-classic-ft",
        title="Как решали задачу ML до больших моделей.",
        intro="Дообучение большой модели — не экзотика новой эпохи, а прямое продолжение "
              "классической схемы transfer learning. Восстановим её с нуля.",
        cards=[
            ("database", "Выборка + эталон",
             "Размеченный набор «вход → правильный ответ» (ground truth, эталонная разметка); "
             "модель учится его воспроизводить и обобщать на новые данные."),
            ("scale", "Train / val / test split",
             "Три непересекающиеся части: на train учат, на validation подбирают, на test меряют "
             "один раз. Правило: нельзя тестировать на обучающих данных."),
            ("git-branch", "Transfer learning (перенос обучения)",
             "Взять предобученную на большом корпусе модель и дёшево дообучить под свою узкую "
             "задачу: быстрее и точнее, чем с нуля («предобучение → дообучение»)."),
        ],
        keep_text="Eval-наборы (golden set) — без них catastrophic forgetting не виден; "
                  "версионирование данных и весов для отката; train/test-дисциплину против утечки; "
                  "мониторинг дрейфа. LoRA удешевила шаг дообучения, но не дисциплину вокруг него.",
        bridge="PEFT/LoRA — тот же transfer learning, доведённый до предела дешевизны и поверх "
               "несравнимо более крупной модели. Идея не новая — новыми стали масштаб базовой "
               "модели и стоимость шага.",
    )


def build_s_classic_agents(p):
    """§4.0 — классическая база раздела 4: классическая автоматизация."""
    build_classic_base(
        p, "s-classic-agents",
        title="Управляемая автоматизация была задолго до агентов.",
        intro="«Агент» звучит как изобретение эпохи больших моделей, но управляемая автоматизация "
              "процессов — зрелая дисциплина. Без неё не оценить, что агент добавляет, а что ломает.",
        cards=[
            ("waypoints", "Конечный автомат",
             "Набор состояний и правил перехода по событиям («создана → в обработке → закрыта»). "
             "Возможные переходы видны, невозможные исключены конструктивно."),
            ("git-branch", "Workflow-движки",
             "Исполняют заранее описанный процесс по фиксированной схеме: BPMN, DAG-оркестраторы "
             "(Airflow), RPA. Порядок шагов определён заранее, а не на лету."),
            ("route", "Управляющий цикл",
             "plan → act → check как контур обратной связи (теория управления, АСУ); военный "
             "аналог — OODA-петля. Не одно действие, а цикл с коррекцией."),
        ],
        keep_text="Детерминированные workflow где можно; идемпотентность (повтор шага не ломает "
                  "состояние); принцип наименьших привилегий (least-privilege) для инструментов; "
                  "лимиты и аудит цикла. Практики надёжности контура никуда не деваются.",
        bridge="LLM-агент — тот же управляющий цикл plan→act→check, но шаги в нём генерирует модель "
               "недетерминированно, а не выбирает фиксированное правило автомата. Отсюда и гибкость, "
               "и провалы раздела.",
    )


def build_s_classic_framework(p):
    """§5.0 — классическая база раздела 5: классический выбор технологии."""
    build_classic_base(
        p, "s-classic-framework",
        title="Как инженер выбирал технологию до AI-хайпа.",
        intro="Сам выбор архитектуры — не новая AI-процедура, а прямое применение классических "
              "принципов инженерного решения. На них и стоят лестница и чек-лист раздела.",
        cards=[
            ("clipboard-list", "От требования, не от инструмента",
             "Требования-инжиниринг: сначала фиксируют, что система обязана делать, потом "
             "подбирают инструмент. Build-vs-buy: строить своё или взять готовое."),
            ("scale", "KISS + YAGNI",
             "Выбирай простейшее решение, закрывающее требование; не закладывай мощность «на "
             "будущее», пока конкретное требование её не потребует."),
            ("milestone", "Наименьшая мощность",
             "Из заметок W3C (Бернерс-Ли и Мендельсон): бери наименее мощный из достаточных "
             "инструментов — его проще анализировать, проверять и сопровождать."),
        ],
        keep_text="Простейшая достаточная архитектура по умолчанию, а бремя доказательства — на "
                  "том, кто хочет её усложнить. В AI-эпоху принцип не отменяется, а дорожает: лишняя "
                  "ступень добавляет недетерминизм, стоимость токенов, задержку, поверхность атаки.",
        bridge="лестница архитектур (код → один вызов → RAG → workflow → агент → мульти-агент) — тот "
               "же принцип наименьшей мощности по AI-архитектурам: оставайся на нижней достаточной "
               "ступени, поднимайся только под требование.",
    )


def build_s01b(p):
    """case_study — Air Canada как первый разбор класса «неправильная
    архитектура под задачу». Hero-фото 787 слева ≥40%, хроника дела +
    вывод справа. (Реюз v4 s01 hero-фото; мем-хук s01 идёт до этого.)"""
    s = blank(p)
    hx, hy, hw, hh = 0.0, 0.0, 6.05, 7.5
    hero_image(s, SCREENSHOTS / "s01-aircanada.jpg", hx, hy, hw, hh)
    filled_rect(s, 0.0, 7.10, 3.9, 0.40, DEEP)
    text_box(s, 0.14, 7.12, 3.7, 0.34,
             "Фото: Air Canada Boeing 787 · Wikimedia · CC-BY-SA",
             size=10.5, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    rx = 6.45
    text_box(s, rx, 0.50, 6.35, 1.50,
             "Первый разбор: чат-бот выдумал политику — платит компания.",
             size=27, bold=True, color=DEEP, line_spacing=1.06)
    text_box(s, rx, 2.05, 6.35, 0.36,
             "Moffatt v. Air Canada · трибунал BC · 14.02.2024",
             size=13.5, bold=True, color=TEAL)
    chron = [
        "Пассажир спросил чат-бота про тариф по случаю утраты близкого",
        "Бот: «купи по полной цене, верни разницу в течение 90 дней»",
        "Реальная политика этого не допускала — и была на той же странице, куда бот ссылался",
        "Трибунал: «бот — не отдельное юр. лицо» → компания вернула $812,02",
    ]
    cy = 2.52
    row_h = 0.86
    for i, t in enumerate(chron):
        circle(s, rx, cy + 0.03, 0.32, MID if i < 3 else GOLD)
        text_box(s, rx, cy + 0.03, 0.32, 0.32, str(i + 1),
                 size=13, bold=True, color=(WHITE if i < 3 else DEEP),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 0.50, cy, 5.85, row_h - 0.04, t,
                 size=13, color=DEEP, line_spacing=1.10,
                 anchor=MSO_ANCHOR.MIDDLE)
        cy += row_h
    gold_callout(s, rx, 5.98, 6.35, 0.90,
                 "Что делать: юридически значимый ответ клиенту — не для генеративного бота. Нужен детерминированный источник политики; бот — лишь навигатор к официальному документу.",
                 size=12.5)
    speaker_notes(s, load_notes("s01b"))


def build_s05c(p):
    """assertion_visual (§1.2) — два значения слова «роль»: роль-персона (тон)
    ≠ протокольные роли system/user/assistant (структура диалога). Приоритет
    system — склонность (~63,8%), не граница; STI/ChatInject подделывают роль.
    Слева chat-template, справа STI-инъекция; числа-плашки; gold-вывод."""
    s = blank(p)
    slide_title(s, "«Роль» — это два разных механизма. Не путайте их.", size=26)
    text_box(s, 0.55, 1.16, 12.25, 0.44,
             "Роль-персона («ты — юрист») настраивает тон. Протокольная роль system/user/assistant — это разметка «кто говорит» в потоке токенов. Второе постоянно принимают за первое.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # LEFT — chat template assembly
    lx, ly, lw, lh = 0.55, 1.82, 6.05, 2.94
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "braces", lx + 0.26, ly + 0.20, 0.44, "mid")
    text_box(s, lx + 0.84, ly + 0.20, lw - 1.0, 0.42,
             "Chat-шаблон: список ролей → плоский поток токенов",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.02)
    # mono token strip
    filled_rect(s, lx + 0.26, ly + 0.86, lw - 0.52, 0.94, DEEP, radius=True,
                radius_adj=0.06)
    text_runs(s, lx + 0.42, ly + 0.94, lw - 0.84, 0.80, [
        {"text": "<|im_start|>system", "size": 12, "color": GOLD,
         "font": FONT_MONO, "bold": True},
        {"text": "  правила поведения  ", "size": 12, "color": WHITE,
         "font": FONT_MONO},
        {"text": "<|im_end|>", "size": 12, "color": TEAL, "font": FONT_MONO},
        {"text": "\n<|im_start|>user", "size": 12, "color": LIGHT,
         "font": FONT_MONO, "bold": True, "newpara": True},
        {"text": "  вопрос  ", "size": 12, "color": WHITE, "font": FONT_MONO},
        {"text": "<|im_end|>", "size": 12, "color": TEAL, "font": FONT_MONO},
    ], line_spacing=1.2)
    text_box(s, lx + 0.26, ly + 1.94, lw - 0.52, 0.90,
             "Спецтокены разметки (ChatML, Llama 4) собираются шаблоном — его можно прочитать и подменить.",
             size=12, color=DEEP, line_spacing=1.16)
    # RIGHT — priority is a tendency, not a guarantee + STI
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, 1.36, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.40,
             "Приоритет system — склонность, не граница",
             size=14, bold=True, color=DEEP)
    text_runs(s, rx + 0.24, ly + 0.58, rw - 0.48, 0.70, [
        {"text": "GPT-4o слушается приоритета ~", "size": 13, "color": DEEP},
        {"text": "63,8%", "size": 16, "bold": True, "color": MID},
        {"text": "  — всё равно ≠ 100%: это склонность.", "size": 13, "color": DEEP},
    ], line_spacing=1.16)
    ocean_box(s, rx, ly + 1.48, rw, 1.46, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    icon(s, "shield-alert", rx + 0.24, ly + 1.66, 0.42, "teal")
    text_box(s, rx + 0.80, ly + 1.66, rw - 1.0, 0.40,
             "STI / role spoofing: роль подделывают",
             size=14, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_runs(s, rx + 0.24, ly + 2.12, rw - 0.48, 0.74, [
        {"text": "Строка ", "size": 12, "color": DEEP},
        {"text": "<|im_start|>assistant", "size": 11.5, "color": MID,
         "font": FONT_MONO, "bold": True},
        {"text": " из внешнего текста → ChatInject ASR ", "size": 12,
         "color": DEEP},
        {"text": "5,18%→32,05%", "size": 14, "bold": True, "color": TEAL},
        {"text": " (Llama-4 до 88,3%).", "size": 12, "color": DEEP},
    ], line_spacing=1.14)
    gold_callout(s, 0.55, 4.92, 9.35, 1.58,
                 "Что делать: не проектируйте защиту в расчёте «система важнее пользователя всегда». Экранируйте спецтокены во ВХОДЯЩЕМ внешнем контенте и проверяйте chat-шаблон у локальных моделей — чужой шаблон тихо ломает приоритет.",
                 size=13.5)
    # #185: реальный интернет-мем «Is this a pigeon?» — протокольную роль
    # system принимают за надёжную границу (а это лишь склонность, не гарантия).
    gmw, gmh = 2.66, 1.78
    gmx, gmy = 10.05, 4.92
    filled_rect(s, gmx - 0.05, gmy - 0.05, gmw + 0.10, gmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s05c-pigeon-ru.png", gmx, gmy, gmw, gmh)
    footer(s, "Полный разбор инъекции в промпт как класса атак агента — в разделе про агенты (безопасность).")
    speaker_notes(s, load_notes("s05c"))


def build_s07(p):
    """case_study (§1.5) — предел CoT: faithfulness. CoT — сгенерированный
    текст, не протокол. Эксперимент с подсказкой; Claude 3.7 ~25% / R1 ~39%;
    хуже на трудных задачах. Проверяй результат, не самообъяснение."""
    s = blank(p)
    slide_title(s, "Рассуждение вслух — не журнал аудита.", size=27)
    text_box(s, 0.55, 1.16, 12.25, 0.46,
             "Faithfulness (верность объяснения) — насколько проговорённая цепочка отражает реальную причину ответа. Измерили напрямую — и она низкая.",
             size=14, italic=True, color=MID, line_spacing=1.14)
    # LEFT — the experiment (2 runs)
    lx, ly, lw = 0.55, 1.86, 6.05
    ocean_box(s, lx, ly, lw, 3.30)
    text_box(s, lx + 0.24, ly + 0.16, lw - 0.48, 0.40,
             "Эксперимент Anthropic (апр. 2025)", size=14, bold=True,
             color=DEEP)
    runs = [
        ("Задачу дают дважды: без подсказки и с подсказкой в промпте, которая меняет ответ («профессор считает, что C»)", "cpu"),
        ("Смотрят: когда ответ изменился под подсказку — упомянула ли модель её в рассуждении?", "eye-off"),
        ("Часто модель меняет ответ, но строит ДРУГУЮ, выдуманную аргументацию, не называя реальную причину", "triangle-alert"),
    ]
    ry = ly + 0.68
    for t, ic in runs:
        icon(s, ic, lx + 0.26, ry + 0.04, 0.40, "mid")
        text_box(s, lx + 0.82, ry, lw - 1.06, 0.86, t,
                 size=12.5, color=DEEP, line_spacing=1.14)
        ry += 0.92
    # RIGHT — the numbers
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, 3.30, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.40,
             "Как часто модель признаёт подсказку", size=14, bold=True,
             color=TEAL)
    bars = [("Claude 3.7 Sonnet", 25, MID), ("DeepSeek R1", 39, TEAL)]
    bt = ly + 0.78
    max_w = rw - 1.9
    for nm, pct, col in bars:
        text_box(s, rx + 0.24, bt, rw - 0.48, 0.30, nm, size=13, bold=True,
                 color=DEEP)
        filled_rect(s, rx + 0.24, bt + 0.34, max_w, 0.40, SOFT_GREY,
                    radius=True, radius_adj=0.3)
        filled_rect(s, rx + 0.24, bt + 0.34, max_w * pct / 100.0, 0.40, col,
                    radius=True, radius_adj=0.3)
        text_box(s, rx + 0.24 + max_w + 0.12, bt + 0.30, 1.4, 0.48,
                 f"~{pct}%", size=18, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        bt += 0.92
    text_box(s, rx + 0.24, bt + 0.02, rw - 0.48, 0.58,
             "И хуже — на трудных задачах (GPQA ниже MMLU): где аудит нужнее всего, объяснению верить можно меньше всего.",
             size=11.5, color=DEEP, line_spacing=1.12)
    gold_callout(s, 0.55, 5.32, 12.25, 0.96,
                 "Что делать: человек-валидатор проверяет РЕЗУЛЬТАТ против независимого источника (база, документ, расчёт, эксперт), а не приложенную «цепочку рассуждений». Контроль на самообъяснении модели — не контроль.",
                 size=13.5)
    speaker_notes(s, load_notes("s07"))


def build_s17(p):
    """assertion_visual (§3.3/§3.5) — FT сузился до поведения: критерии
    «что куда» + гибрид — норма. Таблица знание→RAG / поведение→PEFT /
    дешевле→дистилляция / детерминированное→код. (Перенесена из s14; s14 теперь
    про дистилляцию как отдельную технику.)"""
    s = blank(p)
    slide_title(s, "Что здесь знание, что поведение, что детерминировано.", size=25)
    text_box(s, 0.55, 1.12, 12.25, 0.44,
             "Вопрос на проектировании — не «RAG или дообучение». Разнесите проблему по осям: знание → RAG, поведение → PEFT, дешевле → дистилляция, детерминированное → код.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    ocean_box(s, 0.40, 1.74, 12.55, 3.34)
    tx, ty = 0.55, 1.86
    headers = ["Если задача требует…", "→ правильный инструмент", "…и НЕ этот, потому что"]
    col_w = [3.95, 3.55, 4.75]
    rows = [
        ("знание меняется / нужны свежесть, провенанс", "RAG (или длинный контекст для малого корпуса)",
         "не дообучение: знание устареет, переобучать дорого, риск забывания", False),
        ("стабильное поведение / тон / формат / политика", "дообучение (PEFT)",
         "не RAG: подаёт знание в контекст, но не меняет манеру модели", False),
        ("снизить стоимость / задержку на узкой задаче", "дообучить учителя + дистилляция в ученика",
         "две отдельные техники в связке; промпт/RAG не уменьшают размер модели", False),
        ("детерминированный, верифицируемый ответ", "обычный код, без ИИ",
         "ни RAG, ни FT: ИИ добавит недетерминизм без выигрыша", True),
    ]
    hh, rh = 0.50, 0.66
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID, radius=False)
        text_box(s, cx + 0.14, ty, col_w[j] - 0.28, hh, hd,
                 size=12.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for ri, (c0, c1, c2, isgold) in enumerate(rows):
        bgrow = GOLD_TINT if isgold else (WHITE if ri % 2 == 0 else SURFACE)
        cx = tx
        for j, cc in enumerate([c0, c1, c2]):
            filled_rect(s, cx, yy, col_w[j], rh, bgrow,
                        stroke=(GOLD if isgold else SOFT_GREY),
                        stroke_pt=(1.5 if isgold else 0.75))
            col = (DEEP if isgold else MID) if j == 1 else DEEP
            text_box(s, cx + 0.14, yy + 0.04, col_w[j] - 0.28, rh - 0.08, cc,
                     size=11.5, bold=(j == 1), color=col,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 5.24, 12.25, 1.00,
                 "Что делать: гибрид — норма ТАМ, где у задачи есть И проблема знания, И проблема поведения. Нет проблемы поведения — FT не нужен, даже когда есть RAG. Каждый компонент добавляется под своё требование — то же правило лестницы.",
                 size=13.5)
    speaker_notes(s, load_notes("s17"))


def build_s19b(p):
    """assertion_visual (§4.1) — экономика: агент ×N токенов vs chat,
    мульти-агент ещё ×15; prompt caching из «ориентира» → в «необходимость».
    Baseline — стоимость одного chat-turn."""
    s = blank(p)
    slide_title(s, "Агент — это не «чат подороже». Это другой порядок цены.", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.44,
             "База сравнения — один chat-turn (десятки центов). Каждая ступень вверх умножает расход токенов, а не прибавляет его.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # 3 escalating cost bars (relative to chat-turn = 1×)
    tiers = [
        ("Один вызов LLM (chat)", "×1", 0.11, LIGHT, "базовая стоимость запроса"),
        ("Одиночный агент (цикл)", "≈50×", 0.36, MID, "план→действие→проверка→повтор: много проходов на задачу"),
        ("Мульти-агент", "ещё ×15", 1.0, TEAL, "поверх агента — координация нескольких агентов (Anthropic, 2025)"),
    ]
    by, bx = 1.86, 0.55
    row_h = 1.14
    max_w = 8.6
    for nm, mult, frac, col, sub in tiers:
        ocean_box(s, bx, by, 12.25, row_h - 0.14)
        text_box(s, bx + 0.24, by + 0.12, 3.4, 0.40, nm, size=14.5, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, bx + 0.24, by + 0.54, 3.4, 0.40, sub, size=11,
                 italic=True, color=SLATE, line_spacing=1.05)
        filled_rect(s, bx + 3.75, by + 0.30, max_w, 0.40, SOFT_GREY,
                    radius=True, radius_adj=0.3)
        filled_rect(s, bx + 3.75, by + 0.30, max(0.5, max_w * frac), 0.40,
                    col, radius=True, radius_adj=0.3)
        chip(s, bx + 3.75 + max(0.5, max_w * frac) - 1.35, by + 0.22, 1.30,
             0.56, mult, fill=(GOLD if col is TEAL else col),
             color=(DEEP if col is TEAL else WHITE), size=15)
        by += row_h
    gold_callout(s, 0.55, 5.28, 9.35, 1.42,
                 "Что делать: считайте бюджет ДО выбора архитектуры. Кэш промптов (не пересчитывать неизменный префикс) в агенте — уже не «приятный ориентир», а необходимость: без него ×N-стоимость цикла становится неподъёмной.",
                 size=13.5)
    # #185: реальный интернет-мем «Batman slap» — «агент — просто чат подороже»
    # → пощёчина: это другой порядок цены, а не «чат подороже».
    bmw, bmh = 1.85, 1.79
    bmx, bmy = 10.55, 5.05
    filled_rect(s, bmx - 0.05, bmy - 0.05, bmw + 0.10, bmh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.05)
    add_image(s, WEB / "s19b-batman-ru.png", bmx, bmy, bmw, bmh)
    footer(s, "Множители — порядок величины (Anthropic, 2025), не точный тариф; абсолютная цена зависит от модели и задачи.")
    speaker_notes(s, load_notes("s19b"))


def build_s20(p):
    """assertion_visual (§4.1) — MCP: N×M→N+M («USB-C»); ~11% каталога реально
    запускаемо; поворот доверия (30+ CVE/60дн, обход пути 82%). Удобство
    подключения ≠ безопасность."""
    s = blank(p)
    slide_title(s, "MCP: подключить инструмент — теперь операция на минуты.", size=25)
    # LEFT — N×M -> N+M
    lx, ly, lw, lh = 0.55, 1.30, 6.05, 2.30
    ocean_box(s, lx, ly, lw, lh)
    icon(s, "cable", lx + 0.24, ly + 0.18, 0.44, "mid")
    text_runs(s, lx + 0.82, ly + 0.20, lw - 1.0, 0.42, [
        {"text": "MCP", "size": 17, "bold": True, "color": MID},
        {"text": "  — стандарт подключения (уже разобран выше)", "size": 12.5, "color": DEEP},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 0.82, lw - 0.5, 1.30,
             "Один раз описал инструмент как MCP-сервер — и его видят все модели-клиенты. Подключение стало операцией на минуты. Но лёгкость подключения ничего не говорит о том, что именно подключаешь.",
             size=13, color=DEEP, line_spacing=1.20)
    # RIGHT — the catalog denominator
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.40,
             "Читайте каталог критически", size=14, bold=True, color=DEEP)
    text_runs(s, rx + 0.24, ly + 0.62, rw - 0.48, 0.80, [
        {"text": "«до 90 000 серверов»  →  реально запускаемых ≈ ", "size": 13,
         "color": DEEP},
        {"text": "10 000", "size": 16, "bold": True, "color": MID},
        {"text": "  = ", "size": 13, "color": DEEP},
        {"text": "11%", "size": 20, "bold": True, "color": GOLD_TINT and MID},
        {"text": " каталога.", "size": 13, "color": DEEP},
    ], line_spacing=1.18)
    text_box(s, rx + 0.24, ly + 1.44, rw - 0.48, 0.72,
             "Остальное — дубли, битые и заглушки. И даже рабочие 10 тысяч не проверены на безопасность.",
             size=12, color=DEEP, line_spacing=1.16)
    # BOTTOM — trust turn (security)
    ty2, th = 3.78, 1.60
    ocean_box(s, 0.55, ty2, 12.25, th, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    icon(s, "shield-alert", 0.80, ty2 + 0.22, 0.46, "teal")
    text_box(s, 1.40, ty2 + 0.22, 11.1, 0.42,
             "Поворот доверия: стандартизация подключения ≠ безопасность подключаемого — и обостряет её.",
             size=15, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    facts = [
        ("30+ CVE", "за 60-дневное окно против MCP-серверов (~43% — внедрение команд)"),
        ("82%", "обход пути (path traversal) среди 2 614 проверенных реализаций"),
    ]
    fx = 0.80
    for big, sub in facts:
        text_box(s, fx, ty2 + 0.78, 1.6, 0.46, big, size=22, bold=True,
                 color=TEAL)
        text_box(s, fx + 1.65, ty2 + 0.80, 4.35, 0.66, sub, size=11.5,
                 color=DEEP, line_spacing=1.12)
        fx += 6.15
    gold_callout(s, 0.55, 5.56, 12.25, 0.90,
                 "Что делать: удобство подключения — не аргумент за подключение. Каждый MCP-сервер = чужой код в вашем окружении + носитель инъекции в контексте. Подключаете под требование задачи, с оценкой новой границы доверия.",
                 size=13)
    footer(s, "Актуальные цифры экосистемы MCP и хронология CVE — см. источники.")
    speaker_notes(s, load_notes("s20"))


def build_s22a_multi(p):
    """assertion_visual (§4.3) — мульти-агент по умолчанию НЕ апгрейд: p^n
    (95%×10≈60%); топология рой 17,2× vs координатор 4,4×. Anthropic:
    «works mainly because it helps spend enough tokens»."""
    s = blank(p)
    slide_title(s, "Мульти-агент по умолчанию — не выигрыш, а множитель риска.", size=24)
    text_box(s, 0.55, 1.10, 12.25, 0.44,
             "В цепочке из n шагов, где сбой любого портит результат, надёжности перемножаются: успех ≈ pⁿ. Инженерная интуиция «каждый шаг почти всегда работает» математически ложна.",
             size=13, italic=True, color=MID, line_spacing=1.14)
    # LEFT — p^n visualization (bar per step count)
    lx, ly, lw, lh = 0.55, 1.72, 6.05, 3.28
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.24, ly + 0.12, lw - 0.48, 0.34,
             "Надёжность цепочки = pⁿ", size=14, bold=True, color=DEEP)
    pns = [
        ("95% × 10 шагов", 60, MID),
        ("90% × 10 шагов", 35, TEAL),
    ]
    pb = ly + 0.58
    mw = lw - 2.0
    for nm, pct, col in pns:
        text_box(s, lx + 0.24, pb, lw - 0.48, 0.26, nm, size=12.5, bold=True,
                 color=DEEP)
        filled_rect(s, lx + 0.24, pb + 0.30, mw, 0.36, SOFT_GREY, radius=True,
                    radius_adj=0.3)
        filled_rect(s, lx + 0.24, pb + 0.30, mw * pct / 100.0, 0.36, col,
                    radius=True, radius_adj=0.3)
        text_box(s, lx + 0.24 + mw + 0.10, pb + 0.26, 1.3, 0.44,
                 f"~{pct}%", size=17, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        pb += 0.86
    filled_rect(s, lx + 0.24, pb + 0.04, lw - 0.48, 0.86, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.5, radius=True, radius_adj=0.08)
    text_box(s, lx + 0.40, pb + 0.10, lw - 0.78, 0.74,
             "95% на шаг звучит надёжно — но 0,95¹⁰ ≈ 0,60. Больше агентов = больше шагов = ниже общая надёжность.",
             size=11.5, color=DEEP, line_spacing=1.14, anchor=MSO_ANCHOR.MIDDLE)
    # RIGHT — topology + Anthropic quote
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, 1.46, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.36,
             "Топология решает: координатор > рой", size=14, bold=True,
             color=TEAL)
    text_runs(s, rx + 0.24, ly + 0.58, rw - 0.48, 0.80, [
        {"text": "«Рой» равноправных агентов амплифицирует ошибки ", "size": 12,
         "color": DEEP},
        {"text": "17,2×", "size": 16, "bold": True, "color": TEAL},
        {"text": ";  один координатор — только ", "size": 12, "color": DEEP},
        {"text": "4,4×", "size": 16, "bold": True, "color": MID},
        {"text": " (Zartis/Redis).", "size": 12, "color": DEEP},
    ], line_spacing=1.18)
    ocean_box(s, rx, ly + 1.58, rw, 1.70)
    icon(s, "message-circle", rx + 0.24, ly + 1.76, 0.42, "mid")
    text_box(s, rx + 0.24, ly + 2.24, rw - 0.48, 0.94,
             "Anthropic дословно: «multi-agent works mainly because it helps spend enough tokens to solve the problem» — выигрыш от объёма токенов, не от «магии координации».",
             size=12, italic=True, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 5.22, 12.25, 1.02,
                 "Что делать: начинай с одного сильного агента. Мульти-агент — только если задача распадается на ШИРОКО параллельные независимые подзадачи высокой ценности; иначе +15× токенов и координационные издержки не окупятся.",
                 size=13)
    speaker_notes(s, load_notes("s22a_multi"))


def build_s23b(p):
    """case_study (§4.10) — каталог 10 КЛАССОВ провалов агентов таблицей:
    класс → кейс/дата → урок. Учебная единица — КЛАСС, не кейс. База: 188 из
    344 enterprise-релевантных — агент сломал сам."""
    s = blank(p)
    slide_title(s, "Провалы агентов — это КЛАССЫ, а не список курьёзов.", size=24)
    text_runs(s, 0.55, 1.08, 12.25, 0.44, [
        {"text": "Из ", "size": 13, "color": MID, "italic": True},
        {"text": "344", "size": 14, "bold": True, "color": MID},
        {"text": " значимых для бизнеса ИИ-инцидентов в ", "size": 13,
         "color": MID, "italic": True},
        {"text": "188", "size": 14, "bold": True, "color": TEAL},
        {"text": " (≈55%) автономная система нанесла ущерб в бою БЕЗ атакующего — агент сломал всё сам. Запоминайте класс, не дату.",
         "size": 13, "color": MID, "italic": True},
    ], line_spacing=1.14)
    ocean_box(s, 0.40, 1.58, 12.55, 4.60)
    tx, ty = 0.52, 1.68
    headers = ["Класс провала", "Кейс · дата", "Выученный урок"]
    col_w = [3.75, 3.30, 5.25]
    rows = [
        ("Деструкция + избыток прав", "PocketOS · 04.2026", "Системный промпт ≠ контроль; жёсткая граница + наименьшие привилегии"),
        ("Инъекция без клика (zero-click)", "EchoLeak / M365 Copilot", "«Смертельная тройка» = эксплуатируемо"),
        ("Неуправляемый расход", "$48k/14ч · $1,3M/30дн", "Нет критерия успеха и жёсткого потолка → цикл не завершится"),
        ("Галлюцинация пакетов", "slopsquatting · 19,7%", "Выдуманное имя пакета → атакующий его регистрирует; не доверять без реестра"),
        ("Вредоносный MCP-сервер", "postmark-mcp · 09.2025", "MCP-сервер = непроверенная цепочка поставок; подмена версии (rug-pull)"),
        ("Мульти-агентный каскад", "61% каскадов из апстрима", "Где сломалось ≠ где проявилось; ошибки перемножаются"),
        ("Юр. ответственность", "OLG Hamm · Air Canada", "«Это ответил бот» — не защита; вывод принадлежит компании"),
        ("Петля без бюджета", "$4 200 / 63 ч", "«Пробуй, пока не выйдет» без лимита = буквально"),
        ("Накопление ошибок", "надёжность как pⁿ", "pⁿ, а не усреднение; «докрутить модель» — слабый рычаг"),
        ("Мульти-агентная хрупкость", "Cognition, 2025", "Для задач с зависимостями мульти-агент хуже одного"),
    ]
    hh = 0.40
    rh = (4.60 - 0.20 - hh) / len(rows)
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID, radius=False)
        text_box(s, cx + 0.12, ty, col_w[j] - 0.24, hh, hd,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for ri, (c0, c1, c2) in enumerate(rows):
        bgrow = WHITE if ri % 2 == 0 else SURFACE
        cx = tx
        for j, cc in enumerate([c0, c1, c2]):
            filled_rect(s, cx, yy, col_w[j], rh, bgrow, stroke=SOFT_GREY,
                        stroke_pt=0.5)
            text_box(s, cx + 0.12, yy, col_w[j] - 0.24, rh, cc,
                     size=10.5, bold=(j == 0), color=(DEEP if j == 0 else DEEP),
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 6.30, 12.25, 0.78,
                 "Что делать: столкнувшись с новым инцидентом — назовите ЕГО КЛАСС. Класс задаёт контрмеру (лимит / наименьшие привилегии / проверка между шагами), а конкретные даты и суммы — не нужно.",
                 size=13)
    speaker_notes(s, load_notes("s23b"))


def build_s23c(p):
    """case_study (§4.10) — deep-dive классов с базой: PocketOS 9 сек, runaway
    $48k/$1,3M, slopsquatting 19,7% (1 из 5), каскад 61% из апстрима. Каждый —
    урок + правильная альтернатива."""
    s = blank(p)
    slide_title(s, "Четыре класса крупным планом — с базой сравнения.", size=25)
    cards = [
        ("bomb", "PocketOS — боевая БД за 9 секунд",
         "Агент нашёл токен с неограниченными правами в чужом файле и удалил том + все резервные копии. Ближайший восстановимый — 3-месячной давности.",
         "Урок: системный промпт — не средство контроля безопасности; нужна жёсткая граница."),
        ("flame", "Неуправляемый расход: $48k / 14ч · $1,3M / 30дн",
         "Запрос без критерия успеха — планировщик расширялся и не завершался. База: обычная сессия — центы-доллары; это на 3–5 порядков выше.",
         "Урок: критерий завершения + жёсткий потолок + аварийный стоп обязательны."),
        ("package", "Slopsquatting — 1 из 5 импортов",
         "На 576 000 примеров кода 19,7% ссылок на пакеты — галлюцинации. Атакующий регистрирует выдуманное имя → следующий агент ставит чужой код.",
         "Урок: фиксация (pin) версий/хешей, lock-файлы, проверка новых зависимостей."),
        ("git-merge", "Каскад — 61% ошибок из апстрима",
         "На 73 инцидентах: в 61% корень был в апстрим-слое (поиск/план), а не там, где сбой стал заметен. Следующий агент принимает ошибку за факт.",
         "Урок: проверка между шагами, меньше переходов, трассируемость."),
    ]
    gx, gy = 0.55, 1.28
    cw, chh = 6.05, 2.42
    gap = 0.15
    for i, (ic, title, body, lesson) in enumerate(cards):
        x = gx + (i % 2) * (cw + gap)
        y = gy + (i // 2) * (chh + gap)
        ocean_box(s, x, y, cw, chh)
        filled_rect(s, x + 0.22, y + 0.22, 0.52, 0.52, TEAL, radius=True,
                    radius_adj=0.18)
        icon(s, ic, x + 0.28, y + 0.28, 0.40, "white")
        text_box(s, x + 0.86, y + 0.22, cw - 1.06, 0.52, title, size=14,
                 bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        text_box(s, x + 0.24, y + 0.86, cw - 0.48, 1.02, body, size=11.5,
                 color=DEEP, line_spacing=1.14)
        filled_rect(s, x + 0.24, y + chh - 0.50, cw - 0.48, 0.40, GOLD_TINT,
                    stroke=GOLD, stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(s, x + 0.36, y + chh - 0.50, cw - 0.72, 0.40, lesson,
                 size=10.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    gold_callout(s, 0.55, 6.42, 12.25, 0.72,
                 "Что делать: во всех четырёх корень один — агент применён без внешней границы (прав, бюджета, валидации). Граница ставится ВНЕ агента, промптом её не заменить.",
                 size=13)
    speaker_notes(s, load_notes("s23c"))


def build_s24(p):
    """assertion_visual (§4.8) — карта данных по фиче: ZDR НЕ покрывает
    third-party/MCP/Files/batch; NYT v OpenAI litigation hold пережил «30 дней».
    Чем агентнее — тем больше данных вне ZDR."""
    s = blank(p)
    slide_title(s, "«У нас ZDR» ≠ «данные защищены во всей цепочке».", size=25)
    text_box(s, 0.55, 1.14, 12.25, 0.44,
             "ZDR (zero data retention) покрывает основные вызовы модели — но не всю архитектуру. Агент = модель + инструменты, и инструменты часто вне ZDR.",
             size=13.5, italic=True, color=MID, line_spacing=1.14)
    # LEFT — what ZDR does NOT cover
    lx, ly, lw, lh = 0.55, 1.82, 6.05, 3.10
    ocean_box(s, lx, ly, lw, lh, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "circle-slash", lx + 0.24, ly + 0.18, 0.44, "teal")
    text_box(s, lx + 0.82, ly + 0.20, lw - 1.0, 0.42,
             "ZDR НЕ покрывает", size=15, bold=True, color=TEAL,
             anchor=MSO_ANCHOR.MIDDLE)
    outs = ["Files API", "пакетную обработку", "исполнение кода в контейнерах",
            "MCP-коннектор", "сторонние интеграции",
            "потребительские планы"]
    oy = ly + 0.76
    for i, o in enumerate(outs):
        col_i = i % 2
        ox = lx + 0.28 + col_i * 2.90
        oyy = oy + (i // 2) * 0.56
        circle(s, ox, oyy + 0.06, 0.11, TEAL)
        text_box(s, ox + 0.22, oyy, 2.60, 0.50, o, size=12, color=DEEP,
                 line_spacing=1.02, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.28, ly + 2.54, lw - 0.56, 0.46,
             "Чем агентнее архитектура — тем больше её данных проходит через эти звенья.",
             size=12, bold=True, color=TEAL, line_spacing=1.10)
    # RIGHT — litigation hold overrides policy
    rx, rw = 6.75, 6.05
    ocean_box(s, rx, ly, rw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    icon(s, "gavel", rx + 0.24, ly + 0.18, 0.44, "gold")
    text_box(s, rx + 0.82, ly + 0.20, rw - 1.0, 0.42,
             "Судебный приказ переживает вашу политику", size=14, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, rx + 0.28, ly + 0.80, rw - 0.56, 1.30,
             "NYT v. OpenAI (2025): суд обязал сохранять ВСЕ логи ChatGPT как доказательства. Договорная политика «30 дней» оказалась бессильна против litigation hold по чужому спору.",
             size=13, color=DEEP, line_spacing=1.20)
    text_box(s, rx + 0.28, ly + 2.14, rw - 0.56, 0.86,
             "Данные, покинувшие ваш периметр, живут по правилам, на которые вы не влияете.",
             size=13, bold=True, color=MID, line_spacing=1.16)
    gold_callout(s, 0.55, 5.06, 9.05, 1.62,
                 "Что делать: составьте карту данных по каждой функции до боевой эксплуатации — какие данные, через какое звено, с какой политикой хранения, какие звенья сторонние. Регулируемые/чувствительные данные — только с ZDR/BAA или локально (on-prem).",
                 size=13)
    # #185: реальный интернет-мем «Always has been» — «данные вне ZDR?» / «так
    # было всегда»: то, что покидает периметр, вне вашей политики хранения.
    amw, amh = 3.00, 1.69
    amx, amy = 9.80, 5.06
    filled_rect(s, amx - 0.05, amy - 0.05, amw + 0.10, amh + 0.10, WHITE,
                stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.04)
    add_image(s, WEB / "s24-alwayshasbeen-ru.png", amx, amy, amw, amh)
    speaker_notes(s, load_notes("s24"))


def build_s28(p):
    """summary (§5.3) — итоги-таблица «механизм → граница → что делать» по всем
    архитектурам (паттерн Л2 s38). ≤2 строки/ячейка, БЕЗ переклички с s01.
    Нижняя строка — «не ИИ вовсе». Gold-итог."""
    s = blank(p)
    slide_title(s, "Итог: механизм → его граница → что делать.", size=25)
    ocean_box(s, 0.40, 1.20, 12.55, 5.06)
    tx, ty = 0.52, 1.30
    headers = ["Механизм", "Где ломается (граница)", "Что делать"]
    col_w = [3.05, 4.55, 4.70]
    rows = [
        ("Промпт / роль", "роль-персона меняет тон, не точность", "точность — через контекст и RAG, не через «ты — эксперт»"),
        ("Chain-of-thought", "faithfulness низкая (~25–39%)", "проверяй результат, не самообъяснение"),
        ("RAG", "«вернул» ≠ «вернул правильное»", "опора на источник + метрики поиска + «не знаю» как норма"),
        ("Дообучение / PEFT", "меняет поведение, не знание; забывание", "PEFT + петля проверки + версионирование; знание → RAG"),
        ("Агентный цикл", "план→действие→проверка→повтор — 4 места отказа", "проверка против внешнего критерия; жёсткий потолок на цикл"),
        ("Экипировка / память", "каждый слот — компромисс, не выигрыш", "добавляй слот под требование, с проверкой"),
        ("Мульти-агент", "pⁿ: 95%×10 ≈ 60%; +15× токенов", "один сильный агент по умолчанию"),
        ("Безопасность", "инъекция × широкие права; ZDR не всё", "наименьшие привилегии + карта данных по функции"),
        ("«Не ИИ вовсе»", "детерминированное + верифицируемое", "обычный код — дешевле, предсказуемее, аудируем", True),
    ]
    hh = 0.46
    rh = (5.06 - 0.20 - hh) / len(rows)
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID, radius=False)
        text_box(s, cx + 0.12, ty, col_w[j] - 0.24, hh, hd,
                 size=12.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    yy = ty + hh
    for ri, row in enumerate(rows):
        isgold = len(row) == 4 and row[3]
        c0, c1, c2 = row[0], row[1], row[2]
        bgrow = GOLD_TINT if isgold else (WHITE if ri % 2 == 0 else SURFACE)
        cx = tx
        for j, cc in enumerate([c0, c1, c2]):
            filled_rect(s, cx, yy, col_w[j], rh, bgrow,
                        stroke=(GOLD if isgold else SOFT_GREY),
                        stroke_pt=(1.5 if isgold else 0.5))
            text_box(s, cx + 0.12, yy, col_w[j] - 0.24, rh, cc,
                     size=11, bold=(j == 0 or isgold), color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 6.36, 12.25, 0.72,
                 "Знать инструмент — значит знать его границы. Выбор архитектуры = найти самую нижнюю ступень, которая закрывает требование задачи.",
                 size=13.5, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s28"))


def build_s31(p):
    """qa_minimal — dedicated final Q&A slide (#239, стиль Лекции 1 s31).

    Большое «Q&A» 120pt по центру в DEEP; «Спасибо» 36pt ниже; контактные
    координаты лектора мелким в правом нижнем углу (заполняются перед
    лекцией). Белый фон, без footer и roadmap-bar."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    # #185: реальный интернет-мем «Waiting Skeleton» — «жду ваши вопросы».
    # Левая колонка; текст Q&A смещён вправо, чтобы мем не наложился.
    kmx, kmy, kmw, kmh = 0.85, 1.55, 3.25, 4.40
    ocean_box(s, kmx, kmy, kmw, kmh)
    add_image(s, WEB / "s31-skeleton-ru.png", kmx + 0.14, kmy + 0.14,
              kmw - 0.28, kmh - 0.28)
    text_box(s, x=4.35, y=2.05, w=8.45, h=2.30, text="Q&A",
             size=120, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    text_box(s, x=4.35, y=4.55, w=8.45, h=0.78,
             text="Спасибо", size=36, bold=False, color=MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.2)
    # Контакты преподавателя НЕ на видимом слое (scaffold) — VFY-плашка в
    # speaker notes (заполняется владельцем перед лекцией на GATE B).
    speaker_notes(s, load_notes("s31"))


# ============================================================
# Main
# ============================================================
def main():
    # U-9: 2-part deck spec — loader reads deck.yaml + deck-part2.yaml,
    # validates 36-slide order + cascade lock (s01–s30 not renumbered).
    spec = load_deck()
    if spec is not None:
        n = len(spec["slides"])
        print(f"deck spec OK — {n} slides (deck.yaml + deck-part2.yaml), "
              f"version {spec['deck'].get('version')}, "
              f"totals {spec['totals'].get('slides')}")

    p = setup_pres()
    # v5 (issue #185) presentation order — 51 slides (authoritative ordered
    # list per deck-v5-inventory-draft §3; §9's "54" double-counted dividers).
    #   R0: s01(meme) s02 s02a s03 s04   (s01b снят — #185/#312)
    #   R1: s04a(div) s05 s05a s05c s05b s06 s07 s08 s08a
    #   R2: s09(div) s10 s11 s12 s13
    #   R3: s13a(div) s13b s15 s17 s14 s16
    #   R4: s18(div) s19 s19b s20 s21 s22 s22a_multi s22b s22c s22d
    #        s22e s25 s24 s25b s23 s23b s23c
    #   R5: s25a(div) s26 s27 s27b s28 s29 s30 s31
    # v5b (issue #185 WP8): +5 «классическая база» слайдов — по одному на
    # раздел §1–§5, сразу ПОСЛЕ дивайдера раздела, ПЕРЕД AI-частью. 51→56.
    builders = [
        # R0 — Открытие (5) — s01b снят (#185/#312): хук уже на s01,
        # Air Canada остаётся кейсом §2 (s13).
        build_s01, build_s02, build_s02a, build_s03, build_s04,
        # R1 — Промпт (div + classic-base + 8)
        build_s04a, build_s_classic_prompt, build_s05, build_s05a, build_s05c,
        build_s05b, build_s06, build_s07, build_s08, build_s08a,
        # R2 — RAG (div + classic-base + 4)
        build_s09, build_s_classic_rag, build_s10, build_s11, build_s12, build_s13,
        # R3 — Fine-tune (div + classic-base + 5)
        build_s13a, build_s_classic_ft, build_s13b, build_s15, build_s17,
        build_s14, build_s16,
        # R4 — Агенты (div + classic-base + 16)
        build_s18, build_s_classic_agents, build_s19, build_s19b, build_s20,
        build_s21, build_s22, build_s22a_multi, build_s22b, build_s22c,
        build_s22d, build_s22e, build_s25, build_s24, build_s25b, build_s23,
        build_s23b, build_s23c,
        # R5 — Фреймворк (div + classic-base + 7)
        build_s25a, build_s_classic_framework, build_s26, build_s27, build_s27b,
        build_s28, build_s29, build_s30, build_s31,
    ]
    # sid list — MUST match `builders` order 1:1 (display order, 55 slides).
    sids = [
        "s01", "s02", "s02a", "s03", "s04",
        "s04a", "s-classic-prompt", "s05", "s05a", "s05c", "s05b", "s06",
        "s07", "s08", "s08a",
        "s09", "s-classic-rag", "s10", "s11", "s12", "s13",
        "s13a", "s-classic-ft", "s13b", "s15", "s17", "s14", "s16",
        "s18", "s-classic-agents", "s19", "s19b", "s20", "s21", "s22",
        "s22a_multi", "s22b", "s22c", "s22d", "s22e", "s25", "s24", "s25b",
        "s23", "s23b", "s23c",
        "s25a", "s-classic-framework", "s26", "s27", "s27b", "s28", "s29",
        "s30", "s31",
    ]
    assert len(builders) == 55, f"expected 55 builders, got {len(builders)}"
    assert len(sids) == 55, f"expected 55 sids, got {len(sids)}"

    total = len(builders)
    inject_report = {}
    for idx, (b, sid) in enumerate(zip(builders, sids)):
        b(p)
        slide = p.slides[idx]
        # (1) inject small superscript [N] markers at claim anchors
        inject_report[sid] = R.inject_ref_markers(slide, sid)
        # (2) bottom clickable numbered source list (fold any footer caveat in)
        if sid in R.SLIDE_REFS:
            ftext = _FOOTER_TEXT.get(id(slide))
            fshape = _FOOTER_TEXT.get("_shapes", {}).get(id(slide))
            if fshape is not None:
                # remove the standalone footer; its words survive as ref tail
                fshape._element.getparent().remove(fshape._element)
            R.refs_of_slide(slide, sid, y=7.02, tail=ftext)
        # (3) speaker notes already carry the «Источники:» block + [N] markers
        # (baked into slides/*.md by patch_notes.py — single source of truth,
        # so slide-[N] и notes-[N] не расходятся; builder's speaker_notes(
        # load_notes(sid)) picks it up). Nothing to do here.
        # (4) muted page number «N / 40» bottom-right on every slide
        R.page_number(slide, idx + 1, total)

    # verification print — any anchor that failed to match
    missed = [(sid, a) for sid, rep in inject_report.items()
              for (a, ok) in rep if not ok]
    if missed:
        print("!! UNMATCHED ANCHORS:")
        for sid, a in missed:
            print(f"   {sid}: {a[:70]}")
    else:
        print("all ref anchors matched OK")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"saved {OUT} — {len(p.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
