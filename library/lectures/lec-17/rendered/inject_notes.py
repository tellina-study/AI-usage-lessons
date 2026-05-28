#!/usr/bin/env python3
"""
inject_notes.py — speaker notes injection (Лекция 13 P0 lesson).

Phase 5 для Лекции 17 использует секции `## speaker_notes` (lowercase
underscore) в slides/*.md — НЕ `## Speaker notes`. Этот скрипт адаптирован
под lowercase underscore.

Извлекает полный текст после `## speaker_notes` до следующего `## ` (= `## chapter_ref`)
и инжектирует в PPTX через notes_slide.notes_text_frame.

Запускать ПОСЛЕ build_lec17.py — каждый rebuild стирает notes.

Pipeline:
  python3 render_scatter.py
  python3 build_lec17.py
  python3 inject_notes.py

Issue #145.
"""
import re
import sys
from pathlib import Path

from pptx import Presentation

WT = Path(__file__).resolve().parent.parent     # .../lec-17
PPTX = WT / 'rendered' / 'lec-17.pptx'
SLIDES_DIR = WT / 'slides'

# section dividers + Q&A — short or no notes is acceptable
SHORT_OK = {3, 10, 18, 27, 37}   # 1-based slide numbers (dividers s03/s10/s18/s27 + Q&A s37) — v2 restructure (8c-1)


def extract_speaker_notes(md_path: Path):
    """Извлечь текст после `## speaker_notes` до следующего `## ` (без ###) или EOF."""
    text = md_path.read_text(encoding='utf-8')
    m = re.search(r'^##\s+speaker_notes\s*$', text, re.MULTILINE)
    if not m:
        # fallback to capitalized variant just in case
        m = re.search(r'^##\s+Speaker notes\s*$', text, re.MULTILINE)
    if not m:
        return None
    after = text[m.end():].strip()
    next_h = re.search(r'^##\s+(?!#)', after, re.MULTILINE)
    if next_h:
        after = after[:next_h.start()].strip()
    # collapse hard-wrap newlines inside paragraphs into spaces; keep blank-line paragraph breaks
    paras = re.split(r'\n\s*\n', after)
    paras = [re.sub(r'\s*\n\s*', ' ', para).strip() for para in paras]
    return '\n\n'.join(p for p in paras if p)


def main() -> int:
    if not PPTX.exists():
        print(f'ERROR: PPTX not found: {PPTX}', file=sys.stderr)
        return 1

    p = Presentation(PPTX)
    slide_files = sorted(SLIDES_DIR.glob('s*.md'))
    print(f'Found {len(slide_files)} markdown slides; deck has {len(p.slides)}')

    if len(slide_files) != len(p.slides):
        print('WARNING: count mismatch — proceeding with zip-min', file=sys.stderr)

    injected = 0
    short = []
    for idx, (slide, md) in enumerate(zip(p.slides, slide_files), 1):
        notes_text = extract_speaker_notes(md)
        if notes_text is None:
            print(f's{idx:02d}: NO speaker_notes section in {md.name} — skipping')
            continue
        nf = slide.notes_slide.notes_text_frame
        nf.clear()
        # multi-paragraph
        paras = notes_text.split('\n\n')
        nf.paragraphs[0].text = paras[0]
        for extra in paras[1:]:
            para = nf.add_paragraph()
            para.text = extra
        n = len(notes_text)
        if n < 100 and idx not in SHORT_OK:
            short.append((idx, n))
        print(f's{idx:02d}: injected {n} chars from {md.name}')
        injected += 1

    p.save(PPTX)
    print(f'\nPPTX saved: {injected}/{len(p.slides)} slides got notes injected')
    if short:
        print('WARNING short notes (<100 chars, not divider/Q&A):', short, file=sys.stderr)
    return 0


if __name__ == '__main__':
    sys.exit(main())
