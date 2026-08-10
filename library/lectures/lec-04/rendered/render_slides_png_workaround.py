"""
render_slides_png_workaround.py — per-slide PNG export workaround for
sandboxes where `libreoffice --headless --convert-to pdf` is broken
(see notes/mcp-limitations.md [#162-render-1]: JuNest/proot LibreOffice
AppImage has a missing/incompatible libpdfiumlo.so PDF backend; PNG
export works but only exports slide 1 per invocation).

Isolates each slide into its own single-slide PPTX copy, then LibreOffice
PNG-exports each in turn. Output resolution ~960x720 (lower than the usual
pdftoppm -r 150 path, but sufficient for structural/contrast/overflow QA).

Usage: python3 render_slides_png_workaround.py <deck.pptx> <outdir> [slide_nums_csv]
  slide_nums_csv (optional): e.g. "1,5,12" to render only specific slides.

Set LIBREOFFICE_BIN env var to override the default AppImage AppRun path
if a normal `libreoffice`/`soffice` binary is available instead (prefer
that — this script is a workaround, not the primary render path).
"""
import os
import sys, shutil, subprocess
from pathlib import Path
from pptx import Presentation

def isolate_slide(src_path, slide_idx, out_path):
    p = Presentation(src_path)
    sldIdLst = p.slides._sldIdLst
    ids = list(sldIdLst)
    for i, sldId in enumerate(ids):
        if i != slide_idx:
            sldIdLst.remove(sldId)
    p.save(out_path)

def main():
    src = sys.argv[1]
    outdir = Path(sys.argv[2])
    outdir.mkdir(parents=True, exist_ok=True)
    only = None
    if len(sys.argv) > 3:
        only = set(int(x) for x in sys.argv[3].split(","))
    p = Presentation(src)
    n = len(p.slides._sldIdLst)
    lo = os.environ.get(
        "LIBREOFFICE_BIN",
        "/tmp/claude-999/appimage_extract/squashfs-root/AppRun")
    failed = []
    for i in range(n):
        idx1 = i + 1
        if only is not None and idx1 not in only:
            continue
        tmp_pptx = f"/tmp/_iso_slide_{idx1}.pptx"
        isolate_slide(src, i, tmp_pptx)
        profile = f"/tmp/lo_profile_iso_{idx1}"
        r = subprocess.run(
            ["env", "HOME=/tmp/lo_home5", lo,
             f"-env:UserInstallation=file://{profile}",
             "--headless", "--convert-to", "png", "--outdir", str(outdir),
             tmp_pptx],
            capture_output=True, text=True, timeout=60)
        produced = outdir / f"_iso_slide_{idx1}.png"
        target = outdir / f"s{idx1:02d}.png"
        if produced.exists():
            shutil.move(str(produced), str(target))
            print(f"slide {idx1}: OK")
        else:
            print(f"slide {idx1}: FAILED rc={r.returncode} stderr={r.stderr[-200:]}")
            failed.append(idx1)
    if failed:
        print("FAILED SLIDES:", failed)
        sys.exit(1)

if __name__ == "__main__":
    main()
