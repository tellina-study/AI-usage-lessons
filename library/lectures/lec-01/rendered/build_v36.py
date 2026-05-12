"""
v3.6 build of 6-slide pilot for Lecture 1 (issue #55 v3.6).

Polish round — 7 правок поверх v3.5 (archive-v35/build_v35.py).

Fix 1 — s01 ХУК на открытие: assertion заменён на «Идентификация людей в
        реальном времени — уже с 2023 года на простом ноутбуке.» (вместо
        «Narrow AI работает на ноутбуке без облака — рабочая инженерная
        лошадка, а не магия» — звучало как continuation, не как hook).
        Narrow AI definition остаётся ниже как supporting (italic).
Fix 2 — s02 cover ВИЗУАЛЬНО ОТЛИЧАЕТСЯ от content slides:
        - title 40→64pt;
        - убран Ocean rounded box callout (motif для content, не cover);
        - background tinted SURFACE #F4F7FA (vs WHITE content slides);
        - декоративный «01» в outline 280pt в углу;
        - hero motif перемещён композиционно (правее, крупнее);
        - subtitle одной короткой фразой («Карта применений AI: где
          работает, где — нет.») вместо callout-параграфа.
Fix 3 — s03 убран gold «УГАДАЙ»: восстановлен нейтральный assertion
        «Сначала — ваша оценка, потом — данные.» в navy 28pt bold.
Fix 4 — (то же что Fix 3 в задаче) — purged gold heading.
Fix 5 — Убрать «инженер ИУ6» отовсюду:
        - s02 callout убран целиком (см. Fix 2);
        - s05a subtitle: «преподаватель курса» (без «· ИУ6»);
        - s05b central question: «Где AI работает, где — нет, и как это
          понять?» (вместо «Как инженеру ИУ6 попасть в оставшиеся 10%?»).
Fix 6 — Убраны методические footers:
        - s01 footer удалён (источник YOLOv8 был дублем speaker notes);
        - s02 footer (LO codes + преподаватель) удалён;
        - s03 footer (методический комментарий про reveal) удалён;
        - s04 footer урезан до источников «ВЦИОМ 2025 · Bloomberg 2025»
          (без «По ходу курса заполняем эти слепые зоны»);
        - s05a footer удалён;
        - s05b footer удалён (refs на 14/18/27 и AI-компонента — в notes).
Fix 7 — Тон: убрана претензия на «волшебную пилюлю за 75 минут»:
        - s02 без обещающего callout (Fix 2);
        - s05b central question исследовательский (Fix 5);
        - main takeaway s05b остаётся как есть («Курс — про этот разрыв»
          — это нормально, не recipe).
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

# === Palette ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)  # светло-серо-голубой outline для «01»

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ASSETS = Path("/home/levko/AI-usage-lessons/library/lectures/lec-01/rendered/assets")
OUT = Path("/home/levko/AI-usage-lessons/library/lectures/lec-01/rendered/lec-01-pilot.pptx")
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


# === Helpers ===
def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    """Заполнить фон слайда сплошным цветом (используется для cover tint)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_pt)
    from lxml import etree
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = 0.16
        except Exception:
            pass
    from lxml import etree
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=16, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    from lxml import etree
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def footer(slide, text, *, color=LIGHT, size=12):
    """Footer ТОЛЬКО с источниками. Используется только на s04 в v3.6."""
    text_box(slide, x=0.5, y=7.05, w=12.3, h=0.35, text=text,
             size=size, italic=True, color=color, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)


# ============================================================
# Slide builders
# ============================================================

def build_s01(p):
    """s01 — live_demo. FIX 1: hook на открытие. FIX 6: footer удалён."""
    s = blank(p)
    # FIX 1: новый hook-assertion — звучит как открытие, не как continuation
    text_box(s, x=0.55, y=0.55, w=5.9, h=2.4,
             text="Идентификация людей в реальном времени — уже с 2023 года на простом ноутбуке.",
             size=28, bold=True, color=DEEP, line_spacing=1.18)
    # Narrow AI definition остаётся как supporting (1 строка курсивом)
    text_box(s, x=0.55, y=3.25, w=5.9, h=1.4,
             text=("Narrow AI — модель решает одну задачу "
                   "(обнаружение людей в кадре) и больше ничего."),
             size=15, italic=True, color=MID, line_spacing=1.3)
    # Bottom caption with metric (gold + teal). Two paragraphs.
    tb = s.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(5.9), Inches(1.0))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.line_spacing = 1.35
    for cfg in [
        {"text": "На экране — ", "size": 15, "color": DEEP},
        {"text": "YOLOv8", "size": 15, "color": MID, "bold": True},
        {"text": " на CPU ноутбука: ", "size": 15, "color": DEEP},
        {"text": "31 fps", "size": 15, "color": GOLD, "bold": True},
        {"text": ".", "size": 15, "color": DEEP},
    ]:
        r = p1.add_run(); r.text = cfg["text"]
        r.font.name = FONT_BODY; r.font.size = Pt(cfg.get("size", 15))
        r.font.bold = cfg.get("bold", False); r.font.color.rgb = cfg.get("color", DEEP)
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.35
    for cfg in [
        {"text": "Без интернета", "size": 15, "color": TEAL, "bold": True},
        {"text": "  ·  обучена в 2023.", "size": 15, "color": DEEP},
    ]:
        r = p2.add_run(); r.text = cfg["text"]
        r.font.name = FONT_BODY; r.font.size = Pt(cfg.get("size", 15))
        r.font.bold = cfg.get("bold", False); r.font.color.rgb = cfg.get("color", DEEP)

    # Right column — Ocean rounded box framing the YOLO mock screenshot
    box_x, box_y, box_w, box_h = 6.55, 0.55, 6.3, 4.4
    ocean_box(s, box_x, box_y, box_w, box_h)
    pad = 0.18
    img_w = box_w - 2 * pad
    img_h = img_w * 720.0 / 1280.0
    img_x = box_x + pad
    img_y = box_y + (box_h - img_h) / 2.0
    add_image(s, ASSETS / "illustrations/s01-yolo-mock.png", img_x, img_y, img_w, img_h)

    # Caption под скриншотом — теперь содержит и источник «YOLOv8 (Ultralytics, 2023)»
    # FIX 6: footer был дублем; перенёс минимальное упоминание источника в caption
    text_box(s, x=box_x, y=box_y + box_h + 0.05, w=box_w, h=0.4,
             text="Кадр модели в момент демо: 2 человека в боксах. YOLOv8 (Ultralytics, 2023).",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # FIX 6: НИЖНЕГО footer'а с методическим комментарием больше нет.


def build_s02(p):
    """s02 — DISTINCT cover.
    FIX 2: визуально отличается от content slides:
        - tinted background SURFACE;
        - крупный title 64pt;
        - декоративная «01» в outline ~280pt в правом краю (под hero);
        - убран Ocean rounded box callout;
        - короткий navigational subtitle вместо callout-параграфа;
        - hero motif крупнее, асимметрично.
    FIX 5: убран «инженер ИУ6» (был в callout — callout удалён).
    FIX 6: footer удалён.
    FIX 7: subtitle навигационный («Карта применений AI…»), не обещание.
    """
    s = blank(p)
    # FIX 2: tinted background для cover
    set_slide_bg(s, SURFACE)

    # FIX 2: декоративная «01» — outline-стиль, ОГРОМНАЯ, в правом нижнем углу
    # outline эффект достигаем светло-серо-голубым цветом на tinted фоне
    text_box(s, x=8.0, y=2.7, w=5.3, h=4.7,
             text="01",
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0)

    # Tag: «Лекция» (TEAL chip-style)
    text_box(s, x=0.7, y=1.0, w=6.5, h=0.55,
             text="ЛЕКЦИЯ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    # Underline accent под тегом — короткая teal линия
    filled_rect(s, 0.72, 1.55, 0.7, 0.05, fill=TEAL)

    # FIX 2: крупный title 64pt — single-line wrap должен влезть в 8.0"
    text_box(s, x=0.7, y=2.0, w=8.5, h=2.4,
             text="Введение —\nAI вокруг нас",
             size=64, bold=True, color=DEEP, line_spacing=1.05,
             align=PP_ALIGN.LEFT)

    # FIX 2 + 7: короткий navigational subtitle (НЕ callout-параграф)
    # Маленькая teal линия слева как акцент.
    filled_rect(s, 0.7, 5.45, 0.05, 0.55, fill=TEAL)
    text_box(s, x=0.95, y=5.45, w=8.0, h=0.6,
             text="Карта применений AI: где работает, где — нет.",
             size=22, color=MID, italic=False, align=PP_ALIGN.LEFT,
             line_spacing=1.25)

    # Hero motif — крупнее (5.5" vs 4.6 в v35), композиционно правее-выше
    # Поверх «01» — hero motif доминирует визуально
    hero_w = 5.0
    add_image(s, ASSETS / "illustrations/hero-cover-light.png",
              x=8.0, y=0.9, w=hero_w, h=hero_w)

    # FIX 6: footer (LO codes + преподаватель) УДАЛЁН.
    # Никакого нижнего тёмного текста, минимальная meta снизу не нужна на cover.


def build_s03(p):
    """s03 — poll questions. FIX 3+4: убран gold «УГАДАЙ»; ассертион нейтральный.
    FIX 6: footer удалён."""
    s = blank(p)
    # FIX 3+4: вместо «УГАДАЙ» (gold 40pt) — нейтральный assertion 28pt navy bold
    text_box(s, x=0.6, y=0.55, w=12.2, h=0.85,
             text="Сначала — ваша оценка, потом — данные.",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT,
             line_spacing=1.15)

    card_y = 2.0
    card_h = 4.4
    card_w = 5.95
    pad = 0.25

    # Q1 card
    q1_x = 0.55
    ocean_box(s, q1_x, card_y, card_w, card_h)
    add_image(s, ASSETS / "icons/lucide-hand-blue.png",
              x=q1_x + 0.35, y=card_y + 0.35, w=0.95, h=0.95)
    text_box(s, x=q1_x + 1.5, y=card_y + 0.4, w=card_w - 1.6, h=0.35,
             text="Вопрос 1  ·  выберите ОДИН вариант",
             size=14, bold=True, color=MID, italic=False)
    text_box(s, x=q1_x + 0.4, y=card_y + 1.55, w=card_w - 0.8, h=1.1,
             text="Какой процент россиян использовали AI в 2025?",
             size=22, bold=True, color=DEEP, line_spacing=1.25)

    chip_y = card_y + 3.05
    chip_w = 1.25
    chip_h = 0.55
    gap = 0.13
    options = ["<20%", "20–40%", "40–60%", ">60%"]
    total_w = len(options) * chip_w + (len(options) - 1) * gap
    start_x = q1_x + (card_w - total_w) / 2.0
    for i, opt in enumerate(options):
        chip(s, start_x + i * (chip_w + gap), chip_y, chip_w, chip_h, opt,
             fill=MID, color=WHITE, size=14, bold=True)

    text_box(s, x=q1_x + 0.4, y=chip_y + 0.7, w=card_w - 0.8, h=0.3,
             text="(поднимите руку на одном варианте)",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Q2 card
    q2_x = q1_x + card_w + 0.4
    ocean_box(s, q2_x, card_y, card_w, card_h)
    add_image(s, ASSETS / "icons/lucide-message-square-blue.png",
              x=q2_x + 0.35, y=card_y + 0.35, w=0.95, h=0.95)
    text_box(s, x=q2_x + 1.5, y=card_y + 0.4, w=card_w - 1.6, h=0.35,
             text="Вопрос 2  ·  можно НЕСКОЛЬКО",
             size=14, bold=True, color=TEAL, italic=False)
    text_box(s, x=q2_x + 0.4, y=card_y + 1.55, w=card_w - 0.8, h=1.1,
             text="Кто использовал AI сегодня — и для чего?",
             size=22, bold=True, color=DEEP, line_spacing=1.25)

    chip_y2 = card_y + 3.05
    options2 = ["код", "текст", "перевод", "другое"]
    chip_w2 = 1.4
    total_w2 = len(options2) * chip_w2 + (len(options2) - 1) * gap
    start_x2 = q2_x + (card_w - total_w2) / 2.0
    for i, opt in enumerate(options2):
        chip(s, start_x2 + i * (chip_w2 + gap), chip_y2, chip_w2, chip_h, opt,
             fill=WHITE, stroke=TEAL, color=TEAL, size=14, bold=True)

    text_box(s, x=q2_x + 0.4, y=chip_y2 + 0.7, w=card_w - 0.8, h=0.3,
             text="(поднимите руку столько раз, для скольки задач использовали)",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # FIX 6: footer (методический комментарий про reveal) УДАЛЁН.


def build_s04(p):
    """s04 — poll reveal. FIX 6: footer урезан до источников «ВЦИОМ 2025 · Bloomberg 2025»."""
    s = blank(p)
    text_box(s, x=0.55, y=0.4, w=12.3, h=0.9,
             text="Разница между вашей оценкой и реальностью — карта ваших слепых зон про AI.",
             size=24, bold=True, color=DEEP, line_spacing=1.18)

    box_y = 1.65
    box_h = 4.5

    # Left — donut 51%
    left_x, left_w = 0.55, 5.7
    ocean_box(s, left_x, box_y, left_w, box_h)
    text_box(s, x=left_x + 0.3, y=box_y + 0.2, w=left_w - 0.6, h=0.45,
             text="Проникновение AI в РФ, 2025",
             size=20, bold=True, color=MID, align=PP_ALIGN.CENTER)

    donut_size = 2.65
    donut_x = left_x + (left_w - donut_size) / 2.0
    donut_y = box_y + 0.85
    add_image(s, ASSETS / "charts/c1-vciom-donut.png",
              x=donut_x, y=donut_y, w=donut_size, h=donut_size)

    overlay_h = 0.95
    overlay_y = donut_y + (donut_size - overlay_h) / 2.0 - 0.05
    text_box(s, x=donut_x, y=overlay_y, w=donut_size, h=overlay_h,
             text="51%",
             size=56, bold=True, color=MID, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    text_box(s, x=left_x + 0.3, y=donut_y + donut_size + 0.05, w=left_w - 0.6, h=0.4,
             text="россиян использовали AI",
             size=14, color=DEEP, align=PP_ALIGN.CENTER, bold=True)

    text_box(s, x=left_x + 0.3, y=box_y + box_h - 0.55, w=left_w - 0.6, h=0.4,
             text="ВЦИОМ, 2025  ·  «использовали AI» = генеративные модели",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right — bar chart 4 LLMs
    right_x = left_x + left_w + 0.4
    right_w = 13.333 - 0.55 - right_x
    ocean_box(s, right_x, box_y, right_w, box_h)
    text_box(s, x=right_x + 0.3, y=box_y + 0.2, w=right_w - 0.6, h=0.45,
             text="Использование LLM в РФ, 2025",
             size=20, bold=True, color=MID, align=PP_ALIGN.CENTER)

    bar_w = right_w - 0.4
    bar_h = bar_w * 1100.0 / 1800.0
    if bar_h > box_h - 1.45:
        bar_h = box_h - 1.45
        bar_w = bar_h * 1800.0 / 1100.0
    bar_x = right_x + (right_w - bar_w) / 2.0
    bar_y = box_y + 0.8
    add_image(s, ASSETS / "charts/c2-llm-shares-v35.png",
              x=bar_x, y=bar_y, w=bar_w, h=bar_h)

    text_box(s, x=right_x + 0.3, y=box_y + box_h - 0.55, w=right_w - 0.6, h=0.4,
             text="*Сумма >100% — респонденты могли указать несколько вариантов.*",
             size=13, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # FIX 6: footer урезан до источников БЕЗ заключительной фразы про «по ходу курса заполняем»
    footer(s, "ВЦИОМ 2025  ·  Bloomberg 2025 (DeepSeek share).")


def build_s05a(p):
    """s05a — instructor card.
    FIX 5: убран «· ИУ6» из subtitle — оставлено «преподаватель курса».
    FIX 6: footer (шаблон визитки…) удалён."""
    s = blank(p)
    text_box(s, x=0.55, y=0.55, w=12.3, h=0.9,
             text="Кто я и почему мне это важно.",
             size=28, bold=True, color=DEEP, line_spacing=1.15)

    # Left card — monogram-tile
    left_x, card_y, left_w, card_h = 0.55, 1.7, 4.2, 4.65
    ocean_box(s, left_x, card_y, left_w, card_h)
    mono_size = 2.6
    mono_x = left_x + (left_w - mono_size) / 2.0
    mono_y = card_y + 0.55
    add_image(s, ASSETS / "illustrations/monogram-tile.png",
              x=mono_x, y=mono_y, w=mono_size, h=mono_size)

    # Gold star above monogram
    star_size = 0.4
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(mono_x + mono_size / 2.0 - star_size / 2.0),
                              Inches(card_y + 0.1),
                              Inches(star_size), Inches(star_size))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.fill.background()
    from lxml import etree
    sppr = star._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")

    text_box(s, x=left_x + 0.3, y=card_y + mono_size + 0.85, w=left_w - 0.6, h=0.5,
             text="[Имя Фамилия]",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    # FIX 5: убран « · ИУ6» из subtitle
    text_box(s, x=left_x + 0.3, y=card_y + mono_size + 1.4, w=left_w - 0.6, h=0.4,
             text="преподаватель курса",
             size=14, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right column — 3 motif tiles with icons
    items = [
        ("lucide-briefcase-blue.png", "Опыт с AI", "[N лет]",
         "[работа с моделями / проектами]", TEAL),
        ("lucide-lightbulb-blue.png", "Почему этот курс важен", None,
         "[личная мотивация — заполнит преподаватель]", MID),
        ("lucide-users-blue.png", "Что-то о себе", None,
         "[хобби / факт — снижает дистанцию]", MID),
    ]
    right_x = left_x + left_w + 0.5
    right_w = 13.333 - 0.55 - right_x
    item_h = 1.4
    item_gap = 0.18
    for i, (icon, head, prefix_gold, desc, head_color) in enumerate(items):
        y = card_y + i * (item_h + item_gap)
        ocean_box(s, right_x, y, right_w, item_h)
        add_image(s, ASSETS / "icons" / icon,
                  x=right_x + 0.25, y=y + (item_h - 0.7) / 2.0, w=0.7, h=0.7)
        text_box(s, x=right_x + 1.15, y=y + 0.18, w=right_w - 1.3, h=0.5,
                 text=head, size=20, bold=True, color=head_color)
        if prefix_gold is not None:
            text_runs(s, x=right_x + 1.15, y=y + 0.7, w=right_w - 1.3, h=0.6, runs=[
                {"text": prefix_gold, "size": 14, "bold": True, "color": GOLD},
                {"text": "  ", "size": 14, "color": DEEP},
                {"text": desc, "size": 14, "italic": True, "color": LIGHT},
            ], line_spacing=1.3)
        else:
            text_box(s, x=right_x + 1.15, y=y + 0.7, w=right_w - 1.3, h=0.6,
                     text=desc, size=14, italic=True, color=LIGHT, line_spacing=1.3)

    # FIX 6: footer (шаблон визитки…) УДАЛЁН.


def build_s05b(p):
    """s05b — course frame + central question.
    FIX 5: убран «инженеру ИУ6» из central question.
    FIX 7: central question исследовательский («Где AI работает, где — нет, и как это понять?»).
    FIX 6: footer (refs 14/18/27 + AI-компонента definition) удалён.
    """
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.9,
             text="Главный вопрос курса — не «можно ли AI?», а «НУЖНО ли и ГДЕ?».",
             size=24, bold=True, color=DEEP, line_spacing=1.15)

    # Stake — bold «Стейкс:» + 15pt italic (без последней фразы про «берём верхнюю границу»)
    text_runs(s, x=0.55, y=1.4, w=12.3, h=0.55, runs=[
        {"text": "Стейкс: ", "size": 15, "color": DEEP, "italic": True, "bold": True},
        {"text": "через 3 года ~80% инженерных проектов будут с AI-компонентой ",
         "size": 15, "color": LIGHT, "italic": True},
        {"text": "(Gartner 2025)", "size": 15, "color": LIGHT, "italic": True},
        {"text": ".  В РФ сегодня — ", "size": 15, "color": LIGHT, "italic": True},
        {"text": "только 5–10% доходят до прода ",
         "size": 15, "color": LIGHT, "italic": True, "bold": True},
        {"text": "(АНО ЦЭ 2025).", "size": 15, "color": LIGHT, "italic": True},
    ])

    # Funnel diagram
    funnel_w = 5.9
    funnel_h = funnel_w * 1200.0 / 1500.0
    funnel_x = 0.55
    funnel_y = 2.05
    add_image(s, ASSETS / "diagrams/d2-funnel-90-10-v35.png",
              x=funnel_x, y=funnel_y, w=funnel_w, h=funnel_h)

    # Right column — Ocean box (teal stroke)
    right_x = funnel_x + funnel_w + 0.4
    right_w = 13.333 - 0.55 - right_x
    box_y = 2.05
    box_h = 4.7
    ocean_box(s, right_x, box_y, right_w, box_h, stroke=TEAL)

    # Main takeaway — 24pt bold (оставляем как есть, это не recipe — это констатация разрыва)
    text_runs(s, x=right_x + 0.35, y=box_y + 0.35, w=right_w - 0.7, h=2.6, runs=[
        {"text": "Завтра — ", "size": 24, "color": DEEP, "bold": True},
        {"text": "почти везде", "size": 24, "color": MID, "bold": True},
        {"text": ".\nСегодня — ", "size": 24, "color": DEEP, "bold": True},
        {"text": "почти никто", "size": 24, "color": MID, "bold": True},
        {"text": ".\nКурс — про ", "size": 24, "color": DEEP, "bold": True},
        {"text": "этот разрыв.", "size": 24, "color": GOLD, "bold": True},
    ], line_spacing=1.35)

    div = filled_rect(s, right_x + 0.35, box_y + 2.95, right_w - 0.7, 0.04, fill=MID)

    # FIX 5+7: central question — исследовательский, обезличен
    # Длина строки требует чуть меньшего фонта (24pt) и переноса в 2 строки
    text_runs(s, x=right_x + 0.35, y=box_y + 3.15, w=right_w - 0.7, h=1.4, runs=[
        {"text": "Где AI ", "size": 24, "color": DEEP, "bold": True},
        {"text": "работает", "size": 24, "color": MID, "bold": True},
        {"text": ", где —\n", "size": 24, "color": DEEP, "bold": True},
        {"text": "нет", "size": 24, "color": GOLD, "bold": True},
        {"text": ", и как это понять?", "size": 24, "color": DEEP, "bold": True},
    ], line_spacing=1.3)

    # FIX 6: footer (refs 14/18/27 + AI-компонента) УДАЛЁН.


# ============================================================
# Main
# ============================================================

def main():
    p = setup_pres()
    build_s01(p)
    build_s02(p)
    build_s03(p)
    build_s04(p)
    build_s05a(p)
    build_s05b(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
