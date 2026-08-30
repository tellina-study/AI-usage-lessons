"""
Shared build helpers for Лекция 4 v3 (SDLC re-spine, 37 slides).

Ported from the proven build_lec04.py (32-slide autonomy-ladder axis, issue #99),
adapted for the worktree path + 8-section SDLC roadmap + Russification.

Palette LOCKED: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide. Motif «Ocean rounded box»
(radius 12, surface #F4F7FA, stroke #1C7293 1.5pt) на каждом content-слайде.

Canvas 13.333" × 7.5" (16:9). Fonts fall back to DejaVu (Cyrillic OK).
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x5B, 0x66, 0x78)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFD, 0xF3, 0xDC)
TEAL_TINT = RGBColor(0xE4, 0xF1, 0xF2)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
MID_TINT  = RGBColor(0xE1, 0xEA, 0xF0)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parents[1]      # library/lectures/lec-04
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
CHARTS = ASSETS / "charts-en"
SLIDES_DIR = ROOT / "slides-en"
FONT_HEAD = "DejaVu Sans"
FONT_BODY = "DejaVu Sans"
FONT_MONO = "DejaVu Sans Mono"


# ============================================================
# Canonical URL registry (ПРАВКА 1) — keyed by short id.
# Source: notes/research/lecture-4/references-and-req-engineering.md
# (Deliverable 2). Only URLs from that map are used. Volatile ones
# ([VFY-day-of]) are still linked; the caveat lives in speaker notes.
# ============================================================
URLS = {
    "metr": "https://metr.org/blog/2025-07-10-early-2025-ai-experienced-os-dev-study/",
    "metr_arxiv": "https://arxiv.org/abs/2507.09089",
    "model_spec": "https://model-spec.openai.com/",
    "model_spec_repo": "https://github.com/openai/model_spec",
    "grove": "https://www.youtube.com/watch?v=8rABwKRsec4",
    "fowler_interrogatory": "https://martinfowler.com/bliki/InterrogatoryLLM.html",
    "fowler_harness": "https://martinfowler.com/articles/exploring-gen-ai/harness-engineering-memo.html",
    "fowler_genai": "https://martinfowler.com/articles/exploring-gen-ai.html",
    "fowler_aithoughts": "https://martinfowler.com/articles/202508-ai-thoughts.html",
    "vibe_or_not": "https://martinfowler.com/articles/exploring-gen-ai/to-vibe-or-not-vibe.html",
    "tdd_agent_loop": "https://martinfowler.com/articles/exploring-gen-ai/tdd-in-the-agent-loop.html",
    "nygard_adr": "https://cognitect.com/blog/2011/11/15/documenting-architecture-decisions",
    "adr_org": "https://adr.github.io/",
    "evol_arch": "https://nealford.com/books/buildingevolutionaryarchitectures.html",
    "fitness_fn": "https://www.thoughtworks.com/radar/techniques/architectural-fitness-function",
    "c4": "https://c4model.com/",
    "agents_md": "https://agents.md/",
    "context_rot": "https://www.trychroma.com/research/context-rot",
    "dora_2024": "https://dora.dev/research/2024/dora-report/",
    "dora_2025": "https://dora.dev/dora-report-2025/",
    "osmani_70": "https://addyo.substack.com/p/the-70-problem-hard-truths-about",
    "willison_vibe_eng": "https://simonwillison.net/2025/Oct/7/vibe-engineering/",
    "willison_llms_code": "https://simonwillison.net/2025/Mar/11/using-llms-for-code/",
    "lethal_trifecta": "https://simonwillison.net/2025/Jun/16/the-lethal-trifecta/",
    "brooks": "https://sunnyday.mit.edu/16.355/BrooksNoSilverBullet2.html",
    "beck_tdd": "https://www.informit.com/store/test-driven-development-by-example-9780321146533",
    "stanford_perry": "https://arxiv.org/abs/2211.03622",
    "nyu_asleep": "https://arxiv.org/abs/2108.09293",
    "meta_testgen": "https://arxiv.org/abs/2501.12862",
    "gitclear": "https://www.gitclear.com/ai_assistant_code_quality_2025_research",
    "devin": "https://cognition.ai/blog/swe-bench-technical-report",
    "camoleak": "https://www.legitsecurity.com/blog/camoleak-critical-github-copilot-vulnerability-leaks-private-source-code",
    "slopsquatting": "https://arxiv.org/abs/2406.10279",
    "replit": "https://fortune.com/2025/07/23/ai-coding-tool-replit-wiped-database-called-it-a-catastrophic-failure/",
    "curl": "https://www.theregister.com/2026/01/21/curl_ends_bug_bounty/",
    "anthropic_skill": "https://www.anthropic.com/research/AI-assistance-coding-skills",
    "spec_kit": "https://github.blog/ai-and-ml/generative-ai/spec-driven-development-with-ai-get-started-with-a-new-open-source-toolkit/",
    "spec_kit_repo": "https://github.com/github/spec-kit",
    "kiro": "https://kiro.dev/docs/specs/feature-specs/",
    "ears": "https://research.manchester.ac.uk/en/publications/easy-approach-to-requirements-syntax-ears/",
    "anthropic_playbook": "https://claude.com/blog/the-ai-native-sdlc-playbook",
    "claude_best": "https://code.claude.com/docs/en/best-practices",
    "anthropic_context": "https://www.anthropic.com/engineering/effective-context-engineering-for-ai-agents",
    "tw_complacency": "https://www.thoughtworks.com/radar/techniques/complacency-with-ai-generated-code",
    "fowler_refactoring": "https://martinfowler.com/books/refactoring.html",
    "tw_radar": "https://www.thoughtworks.com/radar",
    # --- added for full [N] coverage (ПРАВКА 1, finish pass) ---
    "willison_llms": "https://simonwillison.net/2025/Mar/11/using-llms-for-code/",
    "dora_google_2025": "https://cloud.google.com/blog/products/ai-machine-learning/announcing-the-2025-dora-report",
    "radar_adr_lw": "https://www.thoughtworks.com/radar/techniques/lightweight-architecture-decision-records",
    "context_rot_repo": "https://github.com/chroma-core/context-rot",
    "spracklen_usenix": "https://www.usenix.org/conference/usenixsecurity25/presentation/spracklen",
    "camoleak": "https://www.legitsecurity.com/blog/camoleak-critical-github-copilot-vulnerability-leaks-private-source-code",
    "cve_59145": "https://nvd.nist.gov/vuln/detail/CVE-2025-59145",
    "register_curl": "https://www.theregister.com/2026/01/21/curl_ends_bug_bounty/",
    "codecrash": "https://arxiv.org/abs/2504.14119",
    "anthropic_skill_arxiv": "https://arxiv.org/abs/2601.20245",
    "adr_templates": "https://github.com/joelparkerhenderson/architecture-decision-record",
    "kiro_specs": "https://kiro.dev/docs/specs/feature-specs/",
    "anthropic_ctx_eng": "https://www.anthropic.com/engineering/effective-context-engineering-for-ai-agents",
    "willison_vibe_code": "https://simonwillison.net/2025/Oct/7/vibe-engineering/",
}


def refs_of(slide, keys, y=6.70, **kw):
    """Convenience: build the bottom numbered ref list from a list of
    (num, name, urlkey) tuples, resolving urlkey via URLS."""
    entries = [(num, name, URLS.get(k, "")) for (num, name, k) in keys]
    return ref_list(slide, entries, y=y, **kw)


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


# ============================================================
# ПРАВКА 2 (owner refinement) — [N] reference markers "существенно
# меньше" основного текста. After a text frame is built, walk its runs,
# find [N] / [N, M] substrings, and re-split them into a smaller (~52%),
# superscript, muted-colour run. Applied automatically by text_box /
# text_runs / gold_callout / teal_callout so the many baked-in [N] markers
# shrink without rewriting every call site. Bottom ref-lists опускают это
# (они и так мелкие).
# ============================================================
_REF_RE = re.compile(r'\[\d+(?:\s*[,–—-]\s*\d+)*\]')
_AMAIN = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _run_props(src_run):
    """Read the font props we need to clone from an existing run."""
    f = src_run.font
    sz = f.size
    return {
        "name": f.name,
        "size_pt": (sz.pt if sz is not None else None),
        "bold": f.bold,
        "italic": f.italic,
        "color": (f.color.rgb if (f.color and f.color.type is not None) else None),
    }


def _clone_run_after(anchor_r, props, text, *, ref=False,
                     ref_frac=0.52, ref_color=LIGHT):
    """Insert a new <a:r> right after anchor_r with cloned props (or a
    small superscript muted variant when ref=True)."""
    new_r = etree.SubElement(anchor_r.getparent(), _AMAIN + "r")
    anchor_r.addnext(new_r)
    rpr = etree.SubElement(new_r, _AMAIN + "rPr")
    base = props["size_pt"] or 16.0
    if ref:
        rpr.set("sz", str(int(round(base * ref_frac * 100))))
        rpr.set("baseline", "30000")
        rpr.set("b", "0")
        rpr.set("i", "1")
    else:
        if props["size_pt"] is not None:
            rpr.set("sz", str(int(round(base * 100))))
        if props["bold"] is not None:
            rpr.set("b", "1" if props["bold"] else "0")
        if props["italic"] is not None:
            rpr.set("i", "1" if props["italic"] else "0")
    # font
    if props["name"]:
        for tag in ("latin", "cs", "ea"):
            el = etree.SubElement(rpr, _AMAIN + tag)
            el.set("typeface", props["name"])
    # colour
    col = ref_color if ref else props["color"]
    if col is not None:
        fill = etree.SubElement(rpr, _AMAIN + "solidFill")
        clr = etree.SubElement(fill, _AMAIN + "srgbClr")
        clr.set("val", str(col))
    t = etree.SubElement(new_r, _AMAIN + "t")
    t.text = text
    return new_r


def shrink_refs_in_frame(text_frame, *, ref_frac=0.52, ref_color=LIGHT):
    """Split every [N] marker inside the frame into a small superscript
    muted run. Non-destructive to surrounding text formatting."""
    for para in text_frame.paragraphs:
        # snapshot runs (we mutate the tree while iterating)
        for run in list(para.runs):
            txt = run.text
            if not txt or "[" not in txt:
                continue
            matches = list(_REF_RE.finditer(txt))
            if not matches:
                continue
            props = _run_props(run)
            # first chunk stays in the original run
            run.text = txt[:matches[0].start()]
            anchor = run._r
            pos = matches[0].start()
            for i, m in enumerate(matches):
                # the marker itself (small)
                anchor = _clone_run_after(anchor, props, m.group(),
                                          ref=True, ref_frac=ref_frac,
                                          ref_color=ref_color)
                # the text between this marker and the next (normal)
                nxt = matches[i + 1].start() if i + 1 < len(matches) else len(txt)
                between = txt[m.end():nxt]
                if between:
                    anchor = _clone_run_after(anchor, props, between, ref=False)
                pos = nxt
    return text_frame


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # support \n as hard paragraph breaks (#sem01-render-1 workaround)
    lines = text.split("\n")
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
    shrink_refs_in_frame(tf)
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before") is not None:
                p.space_before = Pt(cfg["space_before"])
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    shrink_refs_in_frame(tf)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.035, min(0.22, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, fill, *, stroke=None, stroke_pt=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def connector(slide, x1, y1, x2, y2, color=LIGHT, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = etree.SubElement(
            ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
        pd.set("val", dash)
    return cn


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    """[#73-render-1] aspect-safe; [#156-1] fixed h-only branch."""
    path = Path(path)
    if not path.exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path); iw, ih = img.size; img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        ir = iw / ih; br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x),
                                     Inches(y + (h - ah) / 2), width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2),
                                     Inches(y), height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w))
    elif h is not None:                          # [#156-1] fix
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.40, h=0.92, w=12.25, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.10, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 color=DEEP, align=PP_ALIGN.LEFT, stroke_pt=1.5):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=stroke_pt,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.20)


def teal_callout(slide, x, y, w, h, text, *, size=14, bold=False,
                 color=DEEP, align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.18)


def footer(slide, text):
    text_box(slide, x=0.55, y=7.04, w=12.25, h=0.34, text=text,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             line_spacing=1.0)


def src(slide, x, y, w, text, *, size=9, color=LIGHT, align=PP_ALIGN.LEFT,
        h=0.22):
    """Inline muted source caption placed RIGHT AT the material it backs
    (definition / claim / recommendation), not in a bottom footer.
    Small, italic, muted — reads as attribution, not body."""
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, italic=True, color=color, align=align,
             line_spacing=1.0)


def icon(slide, name, x, y, size, variant="mid"):
    add_image(slide, ICONS / f"{name}-{variant}.png", x, y, size, size)


# ============================================================
# ПРАВКА 1 (#269 + #266a) — numbered reference system
# Compact [N] markers at the claim + a small muted CLICKABLE numbered
# source list at the bottom of the slide. URLs come ONLY from the research
# URL map (references-and-req-engineering.md, Deliverable 2).
# ============================================================
def ref_list(slide, entries, *, y=6.70, x=0.55, w=12.25, h=0.60,
             size=8.5, color=LIGHT, line_spacing=1.02, cols=None):
    """Bottom numbered clickable source list.

    entries: list of (num:str, name:str, url:str). Renders «[N] name»
    where name is a clickable hyperlink (run.hyperlink.address = url).
    Muted, italic, small — reads as attribution, never a text-wall.
    Kept to 1–2 visual lines; entries are separated by «   ·   ».
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    for i, (num, name, url) in enumerate(entries):
        # marker [N]
        rm = p.add_run()
        rm.text = f"[{num}] "
        rm.font.name = FONT_BODY; rm.font.size = Pt(size)
        rm.font.bold = True; rm.font.italic = True
        rm.font.color.rgb = MID
        # clickable name
        rn = p.add_run()
        rn.text = name
        rn.font.name = FONT_BODY; rn.font.size = Pt(size)
        rn.font.italic = True
        rn.font.color.rgb = color
        if url:
            try:
                rn.hyperlink.address = url
            except Exception:
                pass
        # separator
        if i < len(entries) - 1:
            rs = p.add_run()
            rs.text = "   ·   "
            rs.font.name = FONT_BODY; rs.font.size = Pt(size)
            rs.font.italic = True
            rs.font.color.rgb = color
    return tb


def link_run(paragraph, text, url, *, size=11, color=MID, bold=False,
             italic=False, font=FONT_BODY):
    """Add a single clickable run to an existing paragraph."""
    r = paragraph.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    if url:
        try:
            r.hyperlink.address = url
        except Exception:
            pass
    return r


def page_number(slide, n, total=None, *, color=SLATE):
    """Small muted page-number stamp in the bottom-right corner.

    Placed at the very bottom-right (x≈12.55, y≈7.16), 10pt italic muted, so it
    never overlaps the left-aligned footer / ref-list (x=0.55) nor the roadmap
    bar (ends y≈7.13). Format «N / TOTAL» when total is given, else «N».
    Applied to every slide by the assembler (build_lec04_v4.py) so all 41
    slides carry it without touching per-slide builders."""
    txt = f"{n} / {total}" if total else str(n)
    tb = slide.shapes.add_textbox(Inches(12.33), Inches(7.16), Inches(0.95),
                                  Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = txt
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = color
    return tb


def speaker_notes(slide, text):
    """Write notes as readable PARAGRAPHS (ПРАВКА 1, owner refinement):
    split on blank lines → one notes-paragraph each, so notes are never a
    single wall of text."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    # normalise: collapse single newlines inside a paragraph to spaces,
    # split into paragraphs on blank lines.
    blocks = [b.strip() for b in re.split(r'\n\s*\n', text.strip()) if b.strip()]
    if not blocks:
        blocks = [""]
    for i, block in enumerate(blocks):
        # keep intentional hard line breaks inside a "Sources:" block
        if block.lstrip().startswith("Sources:"):
            lines = [ln.rstrip() for ln in block.split("\n")]
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = lines[0]
            for ln in lines[1:]:
                sub = tf.add_paragraph()
                sub.text = ln
            continue
        one = re.sub(r'\s*\n\s*', ' ', block)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = one


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# ПРАВКА 1 (owner refinement) — per-slide source registry.
# ONE definition per display-slide id drives BOTH:
#   • the bottom clickable [N] list on the slide (refs_of_slide), and
#   • the «Источники:» block appended to the speaker notes (notes_sources).
# So slide-[N] and notes-[N] can never diverge. Entry:
#   (num:str, short_name:str, urlkey:str, gloss:str[, volatile:bool])
# gloss = one phrase: what the source says / why authoritative.
# volatile → «[VFY-day-of]» appended in notes only.
# URLs resolved ONLY via URLS (research Deliverable 2). No invented URLs.
# ============================================================
SLIDE_REFS = {
    "s01": [
        ("1", "METR — RCT: AI and experienced OSS developers (+19% time)",
         "metr", "randomized controlled trial, n=16, 246 tasks; "
         "measured real time, not perception"),
    ],
    "s04": [
        ("1", "Fowler — Interrogatory LLM", "fowler_interrogatory",
         "the bottleneck of AI development is intent, not typing code"),
        ("2", "Google — DORA 2025", "dora_2025",
         "\"AI amplifies what is already there\": discipline is primary, the tool secondary",
         True),
    ],
    "s05": [
        ("1", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "agentic git loop: each stage commits an artifact, human at the gates"),
        ("2", "OpenAI — Model Spec", "model_spec",
         "spec-as-contract: clause = example prompt = test", True),
        ("3", "GitHub — Spec Kit", "spec_kit",
         "\"intent is the source of truth\"; small verifiable tasks"),
        ("4", "Google — DORA 2025", "dora_2025",
         "seven delivery capabilities; \"AI amplifies what is already there\"", True),
        ("5", "Thoughtworks — Exploring Gen AI (Böckeler)", "fowler_genai",
         "the assistant suggests, the developer owns; harness engineering"),
        ("6", "Willison — Vibe engineering", "willison_vibe_eng",
         "the disciplines an LLM rewards: tests, plans, reviews"),
        ("7", "Brooks — No Silver Bullet", "brooks",
         "essential vs accidental complexity; \"what to build\" is the human's"),
        ("8", "Beck — TDD: By Example", "beck_tdd",
         "red-green-refactor; test-as-specification"),
        ("9", "Nygard — ADR", "nygard_adr",
         "immutable \"why\" records under version control"),
        ("10", "Ford/Parsons — Evolutionary Architectures", "evol_arch",
         "fitness functions: \"fitness\" is objective and automatic"),
        ("11", "Fowler — Refactoring", "fowler_refactoring",
         "the discipline of small verifiable changes"),
        ("12", "Brown — C4", "c4",
         "architecture-as-code: a textual, diffable model for AI"),
    ],
    "s06": [
        ("1", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "each phase commits a human-owned artifact; human at the gates"),
        ("2", "OpenAI — Model Spec", "model_spec",
         "spec-as-contract versioned next to the code", True),
        ("3", "Google — DORA 2025", "dora_2025",
         "seven capabilities independently reproduce the same phase skeleton", True),
        ("4", "Thoughtworks — Exploring Gen AI", "fowler_genai",
         "harness engineering; the assistant suggests, the human owns"),
    ],
    "s07": [
        ("1", "Thoughtworks — Exploring Gen AI (Böckeler)", "fowler_genai",
         "autonomy is a property of the working mode, not the tool brand"),
        ("2", "Willison — using LLMs for code", "willison_llms_code",
         "modes of using an LLM: from autocomplete to orchestrator-agent"),
    ],
    "s08": [
        ("1", "Brooks — No Silver Bullet", "brooks",
         "essential complexity is not removed by a tool — it is in the nature of the task"),
        ("2", "Google — DORA 2025", "dora_2025",
         "\"AI amplifies what is already there\": a multiplier, not a replacement for discipline",
         True),
        ("3", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "the method leads, the tool executes; artifacts and gates live in the practice"),
    ],
    "s10": [
        ("1", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "AI is strong at structuring intent; the intent itself stays with the human"),
        ("2", "OpenAI — Model Spec", "model_spec",
         "requirements are living, versioned, reviewed Markdown, not a prompt",
         True),
        ("3", "GitHub — Spec Kit", "spec_kit",
         "the order requirements -> design -> tasks; small verifiable units"),
        ("4", "Fowler — Interrogatory LLM", "fowler_interrogatory",
         "the bottleneck is expressing intent, not writing code"),
        ("5", "Grove — The New Code (OpenAI)", "grove",
         "\"the source requirements — that is the valuable artifact\"", True),
    ],
    "s11": [
        ("1", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "requirements are reviewed and signed off by a human BEFORE code; accept = \"the merge\""),
        ("3", "GitHub — Spec Kit", "spec_kit",
         "the order requirements -> design -> tasks; versioning next to the code"),
        ("4", "Fowler — Interrogatory LLM", "fowler_interrogatory",
         "the model ASKS questions, surfacing unstated assumptions"),
        ("6", "Thoughtworks — fitness function", "fitness_fn",
         "non-functional requirements are enforced by fitness functions"),
        ("7", "Mavin — EARS notation (IEEE RE'09)", "ears",
         "\"WHEN <trigger>, the system SHALL <response>\" — removes should/may, makes it testable"),
        ("9", "Nygard — ADR", "nygard_adr",
         "keep requirements current like an ADR — in sync with the code"),
    ],
    "s12": [
        ("1", "Brooks — No Silver Bullet", "brooks",
         "the bottleneck is the precision of stating intent (essential complexity)"),
        ("2", "Fowler — Interrogatory LLM", "fowler_interrogatory",
         "the alternative to prompt-and-pray: restore the human checkpoint — requirements before code"),
    ],
    "s14": [
        ("1", "Brooks — No Silver Bullet", "brooks",
         "\"deciding what to build\" is essential complexity, not delegated"),
        ("2", "Thoughtworks — Technology Radar", "tw_radar",
         "codebase cognitive debt — in the Hold ring", True),
        ("3", "Thoughtworks — architectural fitness function", "fitness_fn",
         "the named remedy against erosion — architectural fitness functions"),
    ],
    "s15": [
        ("1", "Nygard — ADR", "nygard_adr",
         "human-written context \"we decided X because Y\" against poisoning"),
        ("2", "Ford/Parsons — Evolutionary Architectures", "evol_arch",
         "fitness functions: deterministic invariants break the loop"),
        ("3", "Brown — C4", "c4",
         "architecture-as-code gives managed, not poisoned, context"),
        ("4", "Thoughtworks — Technology Radar", "tw_radar",
         "a catalog of practices for managing architecture with AI", True),
    ],
    "s16": [
        ("1", "Thoughtworks — Exploring Gen AI (Böckeler)", "fowler_genai",
         "the \"poisoning loop\": AI copies \"how it's done here\", not \"how it's right\""),
        ("2", "Nygard — ADR", "nygard_adr",
         "shared understanding of \"why\" breaks the poisoning loop"),
        ("3", "Thoughtworks — fitness function", "fitness_fn",
         "deterministic invariants against self-entrenchment of bad design"),
    ],
    "s18": [
        ("1", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "the loop explore->plan->code->commit; the order is enforced"),
        ("2", "Osmani — The 70% Problem", "osmani_70",
         "the smaller the AI's proposal, the more real the review; a giant diff goes unread"),
        ("3", "Brooks — No Silver Bullet", "brooks",
         "AI takes the accidental complexity, the human the essential"),
    ],
    "s19": [
        ("1", "agents.md — open standard", "agents_md",
         "a persistent layer in the repository: build/test commands, style, guardrails"),
        ("2", "Chroma — Context Rot (18 models)", "context_rot",
         "retrieval accuracy drops non-linearly as input grows — \"context rots\"",
         True),
        ("3", "Anthropic — context engineering", "anthropic_ctx_eng",
         "3 curation primitives: JIT retrieval, compaction, memory notes"),
    ],
    "s20": [
        ("1", "Thoughtworks — Exploring Gen AI (Böckeler)", "fowler_genai",
         "harness engineering: a deterministic scaffold around the model"),
        ("2", "Willison — Vibe engineering", "willison_vibe_eng",
         "\"review it — or it's not engineering\"; three layers of control"),
        ("3", "Thoughtworks — fitness function", "fitness_fn",
         "invariant violation -> add a fitness function back into the harness"),
    ],
    "s21": [
        ("1", "Osmani — The 70% Problem", "osmani_70",
         "AI speeds up the first ~70%, the last 20-30% (understanding) is hard"),
        ("2", "GitClear — AI code quality 2025 (211M lines)", "gitclear",
         "rising clones and churn, falling refactoring — tech-debt markers (correlation)",
         True),
    ],
    "s22": [
        ("1", "Devin (Cognition) — SWE-bench technical report", "devin",
         "13.86% — only on 25% of the bench, acknowledged contamination, 45-min limit"),
        ("2", "OpenAI — Model Spec", "model_spec",
         "vendor numbers require a comparison baseline and knowledge of the slice", True),
    ],
    "s24": [
        ("1", "Beck — TDD: By Example", "beck_tdd",
         "red-green-refactor; the human owns the test spec"),
        ("2", "Böckeler — TDD in the agent loop", "tdd_agent_loop",
         "TDD-first in the agent loop: no gain + ~3x tokens; the structure is what's valuable",
         True),
        ("3", "Fowler — Exploring Gen AI (tests-as-guardrails)", "fowler_genai",
         "a test forces the interface without coupling to the implementation"),
    ],
    "s25": [
        ("1", "Fowler — Exploring Gen AI", "fowler_genai",
         "\"all green\" lies: the AI's report of a run != proof of a run"),
        ("2", "Meta — TestGen-LLM", "meta_testgen",
         "more coverage (32% vs 5.3%), but fewer mutants killed (2.4% vs 15%)"),
    ],
    "s27": [
        ("1", "Willison — Vibe engineering", "willison_vibe_eng",
         "adversarial review with fresh context; \"review it — or it's not engineering\""),
        ("2", "Osmani — The 70% Problem", "osmani_70",
         "\"if you can't explain it — don't commit\"; accountability is the human's"),
        ("3", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "an adversarial reviewer finds holes even in healthy code — scope on correctness"),
    ],
    "s28": [
        ("1", "Thoughtworks — Complacency with AI-generated code", "tw_complacency",
         "uncritical acceptance of AI code, a drop in critical thinking (Radar Hold)"),
        ("2", "curl — the end of the bug bounty (The Register)", "register_curl",
         "a flood of AI-slop reports: valid share >15% -> <5%; the program was suspended",
         True),
        ("3", "CodeCrash", "codecrash",
         "misleading comments crash the model's reasoning (~-23%)"),
    ],
    "s29": [
        ("1", "Willison — the lethal trifecta", "lethal_trifecta",
         "the lethal trifecta: untrusted content x secrets x egress"),
        ("2", "Fowler — Exploring Gen AI", "fowler_genai",
         "breaking the trifecta is architectural, not \"a better model\""),
        ("3", "Google — Big Sleep / OSS-Fuzz", "dora_google_2025",
         "AI vulnerability hunting — curated cases, not a universal gate", True),
    ],
    "s30": [
        ("1", "Stanford — Perry et al. (CCS 2023)", "stanford_perry",
         "with an AI assistant people introduce vulnerabilities more often and more confidently"),
        ("2", "NYU — Asleep at the Keyboard? (IEEE S&P 2022)", "nyu_asleep",
         "~40% of Copilot programs contained vulnerabilities (in security-sensitive tasks)"),
    ],
    "s31": [
        ("1", "Slopsquatting — Spracklen et al., USENIX Security 2025",
         "spracklen_usenix", "of 576k samples ~20% are non-existent packages; "
         "43% of hallucinations are reproducible"),
        ("2", "CamoLeak — CVE-2025-59145", "camoleak",
         "prompt injection in Copilot Chat -> exfiltration of secrets (CVSS 9.6)"),
    ],
    "s32": [
        ("1", "Fortune — Replit AI wiped the production DB", "replit",
         "the agent broke the code freeze, deleted the prod DB, lied, rated itself 95/100"),
    ],
    "s34": [
        ("1", "Google — DORA 2024", "dora_2024",
         "+throughput and +7.5% documentation, but -7.2% delivery stability", True),
        ("2", "Google — DORA 2025", "dora_2025",
         "a negative link of AI to stability for the second year running", True),
    ],
    "s35": [
        ("1", "Google — DORA 2024", "dora_2024",
         "+7.5% to documentation quality — the only clean plus, but with a paired cost",
         True),
        ("2", "Thoughtworks — Exploring Gen AI (Böckeler)", "fowler_genai",
         "\"AI will not replace a well-documented and automated setup\""),
    ],
    "s37": [
        ("1", "Google — DORA 2025", "dora_2025",
         "\"AI amplifies what is already there\" — the lens the matrix is ordered by",
         True),
        ("2", "Anthropic — AI-Native SDLC playbook", "anthropic_playbook",
         "a phase skeleton with a human-owned artifact at each — the matrix's frame"),
    ],
    "s38": [
        ("1", "Google — DORA (systemic, n~5000)", "dora_2025",
         "throughput+, but AI's link to stability is negative for a second year", True),
        ("2", "GitClear — AI code quality 2025 (211M lines)", "gitclear",
         "refactoring down, duplicates up, churn up — three tech-debt markers (correlation)",
         True),
        ("3", "METR — RCT (n=16, experts)", "metr",
         "+19% time with AI while believing in speed-up — the perception gap"),
    ],
    "s39": [
        ("1", "Böckeler — To vibe or not to vibe", "vibe_or_not",
         "\"using generative AI is a continuous risk assessment\"", True),
    ],
    "s40": [
        ("1", "Anthropic — skill formation (Shen & Tamkin 2026)",
         "anthropic_skill_arxiv", "RCT n=52: quiz 50% with AI vs 67% without (~-17 pp) "
         "when delegating generation; those who asked about concepts show no degradation",
         True),
    ],
}


def _resolve_refs(sid):
    out = []
    for entry in SLIDE_REFS.get(sid, []):
        num, name, urlkey, gloss = entry[0], entry[1], entry[2], entry[3]
        volatile = len(entry) > 4 and entry[4]
        out.append((num, name, URLS.get(urlkey, ""), gloss, volatile))
    return out


def refs_of_slide(slide, sid, *, y=None, size=8.5):
    """Bottom clickable [N] list for a display slide, sourced from
    SLIDE_REFS. Skips silently if the slide has no registry entry."""
    resolved = _resolve_refs(sid)
    if not resolved:
        return None
    entries = [(num, name, url) for (num, name, url, gloss, vol) in resolved]
    yy = y if y is not None else (7.06 if len(entries) <= 2 else 7.02)
    sz = size if len(entries) <= 4 else 8.0
    return ref_list(slide, entries, y=yy, size=sz)


def notes_sources_block(sid):
    """Build the «Источники:» text block for the speaker notes of a display
    slide: numbered [N] + FULL URL + one gloss phrase; volatile → [VFY-day-of].
    Returns "" when the slide has no registry entry."""
    resolved = _resolve_refs(sid)
    if not resolved:
        return ""
    lines = ["Sources:"]
    for (num, name, url, gloss, vol) in resolved:
        vfy = " [VFY-day-of]" if vol else ""
        lines.append(f"[{num}] {name} — {gloss}. {url}{vfy}")
    return "\n".join(lines)


def notes_with_sources(slide, sid):
    """Write speaker notes (paragraph-formatted) with the «Источники:» block
    appended. Single call replaces speaker_notes(slide, load_notes(sid))."""
    body = load_notes(sid)
    block = notes_sources_block(sid)
    text = f"{body}\n\n{block}" if block else body
    speaker_notes(slide, text)


# ============================================================
# Section divider — unified template (8-card roadmap, gold current)
# Sections of Лекции 4 v3 (SDLC): 0..7.
# ============================================================
NAV = [
    ("0", "Intro"),
    ("1", "Requirements"),
    ("2", "Architecture"),
    ("3", "Implementation"),
    ("4", "Testing"),
    ("5", "Review+Sec."),
    ("6", "Delivery+"),
    ("7", "Synthesis"),
]


def roadmap_bar(slide, here_idx, *, y=6.55):
    """8-card progress bar; current section gold-bordered."""
    n = len(NAV)
    gap = 0.10
    bx = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    ch = 0.58
    for i, (num, label) in enumerate(NAV):
        x = bx + i * (cw + gap)
        cur = (i == here_idx)
        if cur:
            filled_rect(slide, x, y, cw, ch, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.12)
        else:
            filled_rect(slide, x, y, cw, ch, SURFACE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(slide, x=x + 0.03, y=y + 0.07, w=cw - 0.06, h=0.20,
                 text=num, size=11, bold=True,
                 color=(DEEP if cur else LIGHT), align=PP_ALIGN.CENTER)
        text_box(slide, x=x + 0.03, y=y + 0.29, w=cw - 0.06, h=0.26,
                 text=label, size=9.5, bold=cur,
                 color=(DEEP if cur else SLATE), align=PP_ALIGN.CENTER,
                 line_spacing=0.95)


def build_section_divider(p, here_idx, subtitle, bridge, sid, tag=None):
    """Distinct divider (NO ocean motif): giant decorative section digit on
    the right (soft outline), РАЗДЕЛ N + subtitle + 1-line narrative bridge on
    the left, gold-current roadmap bar at bottom.

    ПРАВКА 4 (#267): dark phase-tag plate removed from all dividers; only the
    0–7 roadmap navigation bar remains. `tag` kept as ignored kwarg for
    backward-compatible call sites."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=8.55, y=0.20, w=4.5, h=5.8, text=str(here_idx),
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.75, y=1.55, w=7.5, h=0.55,
             text=f"SECTION {here_idx}", size=20, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.18, 0.70, 0.05, fill=GOLD)
    text_box(s, x=0.75, y=2.55, w=7.7, h=1.75, text=subtitle,
             size=32, bold=True, color=DEEP, line_spacing=1.08)
    text_box(s, x=0.78, y=4.45, w=7.65, h=1.85, text=bridge,
             size=15.5, italic=True, color=LIGHT, line_spacing=1.18)
    roadmap_bar(s, here_idx, y=6.55)
    speaker_notes(s, load_notes(sid))
    return s
