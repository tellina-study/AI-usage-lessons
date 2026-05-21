"""
Build script for Лекции 8 «AI в креативных индустриях и медиа» (39 slides).

Phase 5 deck v1 — Media-heavy ≥80% (≥33/39 slides with embedded media).

Source-of-truth: deck.yaml + chapter.md + slides/*.md (39 files with readable
speaker notes 150-300 words).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal
(#028090) secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke
#1C7293 1.5pt).

Canvas: 13.333" × 7.5" (16:9). Pacing ≈ 75 мин.

Build via: python3 build_lec08.py — generates lec-08.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
GREEN_OK  = RGBColor(0x2E, 0x8B, 0x57)
RED_WARN  = RGBColor(0xC0, 0x39, 0x2B)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path("/tmp/lec-08-wt/library/lectures/lec-08")
ASSETS = ROOT / "assets"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-08.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    if not Path(path).exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path)
            img_w, img_h = img.size
            img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        img_ratio = img_w / img_h
        box_ratio = w / h
        if img_ratio > box_ratio:
            actual_h = w / img_ratio
            y_offset = (h - actual_h) / 2
            slide.shapes.add_picture(str(path), Inches(x), Inches(y + y_offset),
                                     width=Inches(w))
        else:
            actual_w = h * img_ratio
            x_offset = (w - actual_w) / 2
            slide.shapes.add_picture(str(path), Inches(x + x_offset), Inches(y),
                                     height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True):
    box = filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                      radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.08, w=w - 0.4, h=h - 0.16, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


def lesson_box(slide, x, y, w, h, text):
    """«УРОК ДЛЯ ИНЖЕНЕРА» gold-tint Ocean rounded box."""
    box = ocean_box(slide, x, y, w, h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.8)
    text_box(slide, x=x + 0.18, y=y + 0.10, w=w - 0.36, h=0.32,
             text="УРОК ДЛЯ ИНЖЕНЕРА",
             size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT, line_spacing=1.0)
    text_box(slide, x=x + 0.18, y=y + 0.45, w=w - 0.36, h=h - 0.55,
             text=text,
             size=14, bold=True, color=DEEP, align=PP_ALIGN.LEFT, line_spacing=1.30)


def speaker_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def load_notes(slide_id):
    """Extract Speaker notes block from slide markdown (дословно copy)."""
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    notes_match = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                             md, re.DOTALL)
    notes = notes_match.group(1).strip() if notes_match else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Section divider — unified template (6 cards)
# ============================================================
NAV_SECTIONS = [
    ("0", "Открытие\n+ ось лекции",     "ось 3 времён\n+ 3 семейства"),
    ("1", "AI ДОБАВИЛ",                 "новые\nвозможности"),
    ("2", "AI ИЗМЕНИЛ",                 "стоимость · скорость\n· профессии"),
    ("3", "AI СЛОМАЛ",                  "12 кейсов\nпровалов"),
    ("4", "AI не нужен",                "4 критерия\nотказа"),
    ("5", "Что делать",                 "5-вопросный\nчек-лист"),
]


def build_section_divider(p, here_idx, title, frame_phrase, notes_slide_id):
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Big section number outline left
    text_box(s, x=0.55, y=1.0, w=4.0, h=4.5,
             text=str(here_idx),
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    # Section label
    text_box(s, x=4.7, y=1.6, w=8.0, h=0.5,
             text="РАЗДЕЛ",
             size=18, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Section title
    text_box(s, x=4.7, y=2.15, w=8.0, h=1.8,
             text=title,
             size=44, bold=True, color=DEEP,
             align=PP_ALIGN.LEFT, line_spacing=1.15)
    # Frame phrase
    filled_rect(s, 4.7, 3.95, 0.04, 0.55, fill=TEAL)
    text_box(s, x=4.85, y=3.95, w=7.9, h=1.0,
             text=frame_phrase,
             size=16, italic=False, color=MID,
             align=PP_ALIGN.LEFT, line_spacing=1.30)
    # Progress bar — 6 cards
    bar_y = 5.85
    n_cells = 6
    total_w = 12.3
    gap = 0.12
    cell_w = (total_w - gap * (n_cells - 1)) / n_cells
    start_x = 0.55
    card_h = 0.95
    for i, (num, sec_title, desc) in enumerate(NAV_SECTIONS):
        x = start_x + i * (cell_w + gap)
        is_here = (i == here_idx)
        if is_here:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=2.0)
            num_color = WHITE
            title_color = WHITE
        elif i < here_idx:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=TEAL_TINT, stroke=TEAL, stroke_pt=1.2)
            num_color = TEAL
            title_color = MID
        else:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=WHITE, stroke=LIGHT, stroke_pt=1.0)
            num_color = LIGHT
            title_color = SLATE
        text_box(s, x=x, y=bar_y + 0.08, w=cell_w, h=0.40, text=num,
                 size=22, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=x + 0.04, y=bar_y + 0.46, w=cell_w - 0.08, h=card_h - 0.50,
                 text=sec_title, size=9, bold=is_here, color=title_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.15)
    speaker_notes(s, load_notes(notes_slide_id))


# ============================================================
# Slide builders — 39 slides
# ============================================================

def build_s01(p):
    """Ice-breaker с hero-image (Sora 2 woolly mammoths) — instant emotional engagement
    + foreshadow keystone «AI добавил → изменил → сломал»."""
    s = blank(p)
    # ---- HERO IMAGE LEFT (Sora 2 mammoths — iconic AI-video frame) ----
    # Image area: 6.6×3.71 (16:9 aspect preserved) ≈ 24.5 sq in. Ocean motif обрамление.
    ocean_box(s, 0.45, 0.55, 6.70, 3.95, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, str(ASSETS / "screenshots/s01-sora-mammoths.png"),
              x=0.55, y=0.65, w=6.50, h=3.75)
    # Attribution chip
    text_box(s, x=0.55, y=4.55, w=6.5, h=0.30,
             text="OpenAI Sora · кадр text-to-video по промпту «woolly mammoths» · 2024",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # ---- ASSERTION + COST-COLLAPSE LEFT-BOTTOM ----
    text_box(s, x=0.55, y=4.95, w=6.6, h=1.05,
             text="AI создаёт артефакт промышленного качества за секунды без специальных навыков.",
             size=20, bold=True, color=DEEP, line_spacing=1.15)
    # Cost-collapse compact 2-row
    text_runs(s, 0.55, 6.10, 6.6, 0.40, [
        {"text": "Музыка: ", "size": 13, "color": DEEP, "bold": True},
        {"text": "композитор + неделя + $500-2000  →  ", "size": 13, "color": SLATE, "italic": True},
        {"text": "Suno · 30 сек · $0", "size": 13, "color": GOLD, "bold": True},
    ], line_spacing=1.25)
    text_runs(s, 0.55, 6.50, 6.6, 0.40, [
        {"text": "Фото: ", "size": 13, "color": DEEP, "bold": True},
        {"text": "фриланс $50-200 + 1-3 дня  →  ", "size": 13, "color": SLATE, "italic": True},
        {"text": "Firefly · 5 сек · $0", "size": 13, "color": GOLD, "bold": True},
    ], line_spacing=1.25)
    # ---- RIGHT: DEMO CARD ----
    ocean_box(s, 7.55, 0.55, 5.30, 5.95)
    text_box(s, x=7.75, y=0.75, w=4.90, h=0.40,
             text="ДЕМО В БРАУЗЕРЕ",
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text_box(s, x=7.75, y=1.20, w=4.90, h=0.50,
             text="Сгенерируем прямо сейчас",
             size=20, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.15)
    # Primary card — Suno
    filled_rect(s, 7.85, 1.95, 4.70, 1.55, MID, radius=True, radius_adj=0.10)
    text_box(s, x=8.00, y=2.05, w=4.40, h=0.32,
             text="ОСНОВНОЕ · АУДИО",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=2.42, w=4.40, h=0.55,
             text="suno.com/create",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=3.00, w=4.40, h=0.50,
             text="тема + жанр + язык → трек 30 сек",
             size=11, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Fallback card — Firefly
    filled_rect(s, 7.85, 3.65, 4.70, 1.55, LIGHT, radius=True, radius_adj=0.10)
    text_box(s, x=8.00, y=3.75, w=4.40, h=0.32,
             text="РЕЗЕРВ · КАРТИНКА",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=4.12, w=4.40, h=0.55,
             text="firefly.adobe.com",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text_box(s, x=8.00, y=4.70, w=4.40, h=0.50,
             text="текст-промпт → фото за 5 сек",
             size=11, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
    # QR placeholder
    filled_rect(s, 11.65, 5.40, 0.85, 0.85, WHITE, stroke=DEEP, stroke_pt=1.5)
    text_box(s, x=11.65, y=5.40, w=0.85, h=0.85, text="QR",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=7.85, y=5.50, w=3.60, h=0.70,
             text="Открой URL на телефоне →",
             size=10, italic=True, color=DEEP, line_spacing=1.20,
             anchor=MSO_ANCHOR.MIDDLE)
    # Footer — Suno / Firefly URLs
    text_box(s, x=0.55, y=7.05, w=12.3, h=0.35,
             text="suno.com · firefly.adobe.com · бесплатный пробный доступ · 30 сек",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Cover slide."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Big outline "08"
    text_box(s, x=8.0, y=2.0, w=5.3, h=5.5, text="08",
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    # Lecture label
    text_box(s, x=0.7, y=1.0, w=6.5, h=0.55, text="ЛЕКЦИЯ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.72, 1.55, 0.7, 0.05, fill=TEAL)
    # Main title
    text_box(s, x=0.7, y=2.0, w=8.5, h=2.6,
             text="AI в творческих\nиндустриях и медиа",
             size=54, bold=True, color=DEEP, line_spacing=1.05,
             align=PP_ALIGN.LEFT)
    # Subtitle
    filled_rect(s, 0.7, 5.45, 0.05, 0.6, fill=TEAL)
    text_box(s, x=0.95, y=5.45, w=10.5, h=0.7,
             text="Что AI добавил, изменил и сломал — и где сказать «нет».",
             size=20, color=MID, italic=False, align=PP_ALIGN.LEFT,
             line_spacing=1.25)
    # Gold highlight chip — student-facing tagline
    chip(s, 0.7, 6.5, 3.6, 0.45, "Лекция 08 · 75 мин · 39 слайдов",
         fill=GOLD_TINT, stroke=GOLD, color=DEEP, size=11, bold=True)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    """Central question."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=0.5, w=12.3, h=0.7,
             text="ЦЕНТРАЛЬНЫЙ ВОПРОС ЛЕКЦИИ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
             line_spacing=1.0)
    # Big Ocean box with the question
    ocean_box(s, 1.2, 1.5, 10.9, 4.5)
    text_runs(s, 1.5, 1.8, 10.3, 4.0, [
        {"text": "Что AI сделал ", "size": 36, "color": GOLD, "bold": True},
        {"text": "с креативной индустрией к 2026 —", "size": 36, "color": DEEP, "bold": True},
        {"newpara": True, "text": "и где инженеру разумно сказать", "size": 36, "color": DEEP, "bold": True},
        {"newpara": True, "text": "«здесь AI не нужен»?", "size": 36, "color": TEAL, "bold": True},
    ], line_spacing=1.30, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Anchor below box
    text_box(s, x=0.55, y=6.3, w=12.3, h=0.5,
             text="Двусоставен по проекту — обе части равно важны.",
             size=15, italic=True, color=MID, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """Lecture map — 6-card horizontal дорожная карта."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.7,
             text="Карта лекции — 6 разделов + Q&A",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="6 разделов · 75 мин · 1 чек-лист на выходе",
             size=16, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 6-card horizontal дорожная карта
    bar_y = 2.2
    n_cells = 6
    total_w = 12.3
    gap = 0.18
    cell_w = (total_w - gap * (n_cells - 1)) / n_cells
    start_x = 0.55
    card_h = 4.2
    here_idx = 0  # Currently in Section 0
    for i, (num, sec_title, desc) in enumerate(NAV_SECTIONS):
        x = start_x + i * (cell_w + gap)
        is_here = (i == here_idx)
        if is_here:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=2.0)
            num_color = WHITE
            title_color = WHITE
            desc_color = WHITE
        else:
            ocean_box(s, x, bar_y, cell_w, card_h,
                      fill=WHITE, stroke=LIGHT, stroke_pt=1.2)
            num_color = LIGHT
            title_color = DEEP
            desc_color = SLATE
        text_box(s, x=x, y=bar_y + 0.3, w=cell_w, h=1.5, text=num,
                 size=72, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=x + 0.10, y=bar_y + 2.0, w=cell_w - 0.20, h=1.0,
                 text=sec_title, size=14, bold=True, color=title_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.20)
        text_box(s, x=x + 0.10, y=bar_y + 3.0, w=cell_w - 0.20, h=1.0,
                 text=desc, size=10, color=desc_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)
    # Q&A separately
    chip(s, 5.5, 6.7, 2.5, 0.5, "+ Q&A",
         fill=TEAL_TINT, stroke=TEAL, color=DEEP, size=14, bold=True)
    speaker_notes(s, load_notes("s04"))


def build_s05(p):
    """KEYSTONE — AI добавил → изменил → сломал."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Main title
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.7,
             text="AI добавил → изменил → сломал",
             size=40, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.25, w=12.3, h=0.5,
             text="Три времени одного процесса — каждое поколение творческих инструментов проходит их за месяцы.",
             size=16, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 timing bands — equal width, stacked vertically
    band_x = 0.55
    band_w = 12.3
    band_y_start = 2.1
    band_h = 1.45
    band_gap = 0.15
    bands = [
        ("ДОБАВИЛ", "Новые возможности — то, чего раньше технологически не было.",
         "Sora 2 · 25 сек text-to-video · sync audio. Midjourney Omni Reference. ElevenLabs — клон голоса из 1 мин.",
         MID, WHITE),
        ("ИЗМЕНИЛ", "Стоимость · скорость · профессии — новая экономика креативных индустрий.",
         "Обвал стоимости 100×–10 000× по типам ассетов. Концепт-арт: дни → секунды. Upwork +70% год к году в AI/ML.",
         LIGHT, WHITE),
        ("СЛОМАЛ", "Новый класс провалов и юридического долга.",
         "RIAA против Suno, миллионы Sony. Arup CFO дипфейк $25.6 млн. SI — фейковые авторы. NYT против OpenAI — 20 млн логов.",
         GOLD, DEEP),
    ]
    for i, (label, subtitle, example, fill_color, txt_color) in enumerate(bands):
        y = band_y_start + i * (band_h + band_gap)
        # Color band
        filled_rect(s, band_x, y, 2.5, band_h, fill_color, radius=True, radius_adj=0.10)
        text_box(s, x=band_x + 0.20, y=y + 0.10, w=2.2, h=0.5, text=str(i+1),
                 size=22, bold=True, color=txt_color, align=PP_ALIGN.LEFT,
                 line_spacing=1.0)
        text_box(s, x=band_x + 0.20, y=y + 0.60, w=2.2, h=0.7,
                 text=label, size=22, bold=True, color=txt_color,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        # Subtitle + example
        ocean_box(s, band_x + 2.7, y, band_w - 2.7, band_h)
        text_box(s, x=band_x + 2.85, y=y + 0.15, w=band_w - 2.95, h=0.5,
                 text=subtitle, size=15, bold=True, color=DEEP,
                 align=PP_ALIGN.LEFT, line_spacing=1.20)
        text_box(s, x=band_x + 2.85, y=y + 0.65, w=band_w - 2.95, h=0.75,
                 text=example, size=12, italic=True, color=SLATE,
                 align=PP_ALIGN.LEFT, line_spacing=1.30)
    # Anchor footer
    text_box(s, x=0.55, y=7.0, w=12.3, h=0.35,
             text="Это ось трёх времён одного процесса, не три параллельные категории.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s05"))


def build_s05a(p):
    """3 families: diffusion / latent video transformer / neural audio."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.7,
             text="3 семейства генеративных моделей медиа",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Ментальная модель: каждое семейство → свои фундаментальные ограничения, не зависящие от качества реализации.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 cards equal width
    card_y = 2.0
    card_h = 4.7
    total_w = 12.3
    gap = 0.20
    card_w = (total_w - gap * 2) / 3
    start_x = 0.55
    families = [
        ("1. Диффузионные",
         "шум → обращение → изображение",
         "Stable Diffusion · Midjourney · Flux · DALL-E · Imagen · Adobe Firefly",
         "«Коммерческая безопасность» зависит от обучающего корпуса, а не от архитектуры.",
         MID),
        ("2. Латентные видео-трансформеры",
         "латентное пространство + временная консистентность",
         "Sora 2 · Veo 3.1 · Runway · Kling 3.0",
         "25-сек предел Sora — стоимость растёт линейно + временной дрейф после ~25 сек.",
         LIGHT),
        ("3. Нейросинтез аудио",
         "авторегрессионный + диффузия",
         "Suno · Udio · ElevenLabs · Stable Audio",
         "Клон голоса из 1 мин — дообучение предобученной foundation-модели, а не с нуля.",
         GOLD),
    ]
    for i, (title, principle, tools, consequence, accent) in enumerate(families):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        # Header band
        filled_rect(s, x, card_y, card_w, 0.55, accent,
                    radius=True, radius_adj=0.10)
        title_color = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=card_y + 0.10, w=card_w - 0.30, h=0.40,
                 text=title, size=17, bold=True, color=title_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Principle
        text_box(s, x=x + 0.18, y=card_y + 0.75, w=card_w - 0.36, h=0.4,
                 text="ПРИНЦИП", size=10, bold=True, color=TEAL,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 1.05, w=card_w - 0.36, h=0.5,
                 text=principle, size=14, bold=True, color=DEEP, italic=True,
                 align=PP_ALIGN.LEFT, line_spacing=1.25)
        # Tools
        text_box(s, x=x + 0.18, y=card_y + 1.75, w=card_w - 0.36, h=0.4,
                 text="ИНСТРУМЕНТЫ 2026", size=10, bold=True, color=TEAL,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 2.05, w=card_w - 0.36, h=1.1,
                 text=tools, size=12, color=SLATE,
                 align=PP_ALIGN.LEFT, line_spacing=1.40)
        # Consequence
        text_box(s, x=x + 0.18, y=card_y + 3.25, w=card_w - 0.36, h=0.4,
                 text="ИНЖЕНЕРНОЕ СЛЕДСТВИЕ", size=10, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 3.55, w=card_w - 0.36, h=1.05,
                 text=consequence, size=12, color=DEEP, bold=True,
                 align=PP_ALIGN.LEFT, line_spacing=1.30)
    speaker_notes(s, load_notes("s05a"))


# Section dividers handled by build_section_divider helper
def build_s06(p):
    build_section_divider(p, here_idx=1, title="AI ДОБАВИЛ",
                          frame_phrase="Новые возможности: text-to-video · сохранение персонажа между генерациями · клонирование голоса · модели мира.",
                          notes_slide_id="s06")


def build_s07(p):
    """Text-to-video 2026 — 3-card comparison."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Text-to-video 2026 — промышленного качества",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="3 флагманские модели определяют состояние индустрии",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3-card comparison
    card_y = 2.0
    card_h = 4.0
    total_w = 8.0
    gap = 0.18
    card_w = (total_w - gap * 2) / 3
    start_x = 0.55
    models = [
        ("Sora 2", "OpenAI", "25 сек · 1080p · sync audio",
         "$0.10/сек 720p · cameos", MID, "openai.com/index/sora-2/"),
        ("Veo 3.1", "Google", "8 сек · 720p/1080p · native audio",
         "$0.05/сек Lite", LIGHT, "Google AI Ultra"),
        ("Kling 3.0", "Kuaishou", "15 сек · 4K · 60 fps",
         "ELO #1 (1243)", GOLD, "5 фев 2026 · 60M creators"),
    ]
    for i, (name, vendor, specs, price, accent, url) in enumerate(models):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        # Accent header
        filled_rect(s, x, card_y, card_w, 0.5, accent, radius=True, radius_adj=0.10)
        name_color = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.12, y=card_y + 0.06, w=card_w - 0.24, h=0.4,
                 text=name, size=18, bold=True, color=name_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.18, y=card_y + 0.65, w=card_w - 0.36, h=0.4,
                 text=vendor, size=11, bold=True, color=TEAL, italic=True)
        text_box(s, x=x + 0.18, y=card_y + 1.10, w=card_w - 0.36, h=1.0,
                 text=specs, size=13, color=DEEP, bold=True, line_spacing=1.35)
        text_box(s, x=x + 0.18, y=card_y + 2.20, w=card_w - 0.36, h=1.0,
                 text=price, size=13, color=GOLD, bold=True, line_spacing=1.30)
        text_box(s, x=x + 0.18, y=card_y + 3.20, w=card_w - 0.36, h=0.7,
                 text=url, size=10, italic=True, color=SLATE, line_spacing=1.20)
    # Right side: real Sora кадр из демо-релиза (woolly mammoths walking through snow,
    # OpenAI Sora YouTube channel · эталонная demo from launch)
    ocean_box(s, 8.75, 2.0, 4.20, 4.0)
    text_box(s, x=8.95, y=2.10, w=3.8, h=0.35,
             text="OPENAI SORA · ДЕМО ПРИ ЗАПУСКЕ · YOUTUBE",
             size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_image(s, str(ASSETS / "screenshots/s07-sora2-frame.png"),
              x=8.92, y=2.50, w=3.86, h=3.20)
    text_box(s, x=8.95, y=5.75, w=3.8, h=0.30,
             text="openai.com/index/sora-2/  ·  $0.10/сек",
             size=10, italic=True, color=LIGHT, line_spacing=1.15)
    # Мини-провал block
    lesson_box(s, 0.55, 6.2, 12.3, 1.0,
               "Мини-провал: 25 сек — это не фильм. Кинематографический пайплайн собирается из коротких блоков с человеческим руководством (Lionsgate × Runway = усиление, а не замещение).")
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """Сохранение персонажа."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Сохранение персонажа: cameos и Omni Reference",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Нарратив из нескольких сцен стал возможен. Без этого AI-видео не работает для истории.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 cards vertical (left) — compressed to fit alongside lesson box
    cards_x = 0.55
    cards_w = 6.0
    card_h = 1.18
    card_gap = 0.10
    cards = [
        ("Sora 2 cameos", "Партнёрство OpenAI + Disney на $1+ млрд",
         "Регистрируется персонаж — вызывается по имени в промптах из нескольких сцен."),
        ("Midjourney Omni Reference (v7)",
         "Точность image-to-image 60% → 85%+",
         "Референс-изображение персонажа сохраняет пропорции лица/одежды/посадки."),
        ("Runway Gen-4 Director Mode",
         "Скриптинг между сценами через структурированные объекты",
         "Персонаж + локации + паттерны движения как константы для всех сцен."),
    ]
    for i, (title, subtitle, body) in enumerate(cards):
        y = 2.0 + i * (card_h + card_gap)
        ocean_box(s, cards_x, y, cards_w, card_h)
        text_box(s, x=cards_x + 0.18, y=y + 0.08, w=cards_w - 0.36, h=0.32,
                 text=title, size=14, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=cards_x + 0.18, y=y + 0.40, w=cards_w - 0.36, h=0.28,
                 text=subtitle, size=10, italic=True, color=TEAL, line_spacing=1.0)
        text_box(s, x=cards_x + 0.18, y=y + 0.68, w=cards_w - 0.36, h=0.48,
                 text=body, size=11, color=SLATE, line_spacing=1.20)
    # Right: real Midjourney референс персонажа grid («рыцарь в лесу» + «старик» с шляпой) —
    # 8 generated images showing character preservation across cw=0 vs cw=100 (consistency weight)
    ocean_box(s, 7.0, 2.0, 5.95, 3.95)
    text_box(s, x=7.15, y=2.10, w=5.6, h=0.35,
             text="MIDJOURNEY · РЕФЕРЕНС ПЕРСОНАЖА · «рыцарь в лесу» × «старик»",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s08-character-grid.png"),
              x=7.10, y=2.45, w=5.78, h=3.40)
    # Lesson box — positioned just below 3-card stack & grid
    lesson_box(s, 0.55, 6.05, 12.3, 1.20,
               "Дрейф между сценами проявляется после 5-10 сцен. Континьюити-супервайзер — новая профессия в творческом пайплайне.")
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """Клонирование голоса + дубляж на несколько языков."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Клонирование голоса + дубляж на несколько языков",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="ElevenLabs: клон голоса из 1 мин аудио → 32+ языков. Промышленное применение в корпоративном секторе.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: real ElevenLabs официальная обложка (from elevenlabs.io/cover.png)
    ocean_box(s, 0.55, 2.0, 6.0, 4.0)
    text_box(s, x=0.75, y=2.10, w=5.6, h=0.35,
             text="ELEVENLABS · ОФИЦИАЛЬНО · elevenlabs.io",
             size=12, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_image(s, str(ASSETS / "screenshots/s09-elevenlabs.png"),
              x=0.75, y=2.45, w=5.60, h=2.10)
    # Список голосов — concise, soundalike risk shown as Teal teal
    voices = [
        ("Голос 1", "Многоязычный · выразительный", MID),
        ("Голос 2", "Russian · neutral", MID),
        ("Голос 3", "EN-UK · narrator", LIGHT),
        ("Голос 4", "ScarJo-like soundalike", TEAL),
    ]
    for i, (vname, desc, badge_c) in enumerate(voices):
        y = 4.65 + i * 0.32
        filled_rect(s, 0.75, y, 5.60, 0.28, SURFACE, stroke=LIGHT, stroke_pt=0.5,
                    radius=True, radius_adj=0.20)
        filled_rect(s, 0.82, y + 0.05, 0.18, 0.18, badge_c, radius=True, radius_adj=0.5)
        text_box(s, x=1.10, y=y, w=5.20, h=0.28,
                 text=f"{vname}  ·  {desc}",
                 size=10, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # Right top: обвал стоимости arrow
    ocean_box(s, 6.85, 2.0, 6.10, 1.8)
    text_box(s, x=7.05, y=2.15, w=5.7, h=0.4,
             text="ОБВАЛ СТОИМОСТИ: дубляж на язык",
             size=12, bold=True, color=TEAL)
    text_runs(s, 7.05, 2.6, 5.7, 0.8, [
        {"text": "$50-500", "size": 24, "color": SLATE, "bold": True, "italic": True},
        {"text": "  →  ", "size": 26, "color": DEEP, "bold": True},
        {"text": "<$1", "size": 32, "color": GOLD, "bold": True},
        {"text": "  /мин", "size": 18, "color": DEEP, "bold": True},
    ], line_spacing=1.10)
    text_box(s, x=7.05, y=3.40, w=5.7, h=0.35,
             text="Длинный формат: недели → минуты (Dubbing Studio, 29 языков)",
             size=11, italic=True, color=LIGHT)
    # Right middle: productionrs
    text_box(s, x=6.85, y=3.95, w=6.10, h=0.4,
             text="ПРОМЫШЛЕННОЕ ПРИМЕНЕНИЕ",
             size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    chip(s, 6.85, 4.35, 1.95, 0.45, "Deutsche Telekom",
         fill=MID, color=WHITE, size=12)
    chip(s, 8.95, 4.35, 1.85, 0.45, "Klarna",
         fill=MID, color=WHITE, size=12)
    chip(s, 10.95, 4.35, 2.00, 0.45, "Multi-language",
         fill=LIGHT, color=WHITE, size=11)
    # Мини-провал ScarJo
    lesson_box(s, 6.85, 5.05, 6.10, 1.95,
               "МИНИ-ПРОВАЛ: ScarJo против OpenAI «Sky» (май 2024). OpenAI убрал голос за неделю — формально без иска. De-facto win для права на образ: клон голоса обязывает к явное согласие, даже если технологически «всего лишь похож».")
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    """Genie 3 world models."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Модели мира — Genie 3 (DeepMind)",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Текст → играбельный 3D-мир @ 24 fps. Это НЕ видеогенерация — это симулированное окружение.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: real Genie 3 official 9-frame gameplay grid from DeepMind blog
    # (volcano, jellyfish, eagle, Japan, waterfall, Venice, wingsuit, alley + WASD controls)
    ocean_box(s, 0.55, 2.0, 7.0, 4.5)
    text_box(s, x=0.75, y=2.10, w=6.6, h=0.35,
             text="GENIE 3 · DEEPMIND · ОФИЦИАЛЬНОЕ ДЕМО",
             size=12, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s10-genie3-world.png"),
              x=0.70, y=2.50, w=6.70, h=3.55)
    text_box(s, x=0.75, y=6.10, w=6.6, h=0.30,
             text="deepmind.google/blog/genie-3-a-new-frontier-for-world-models",
             size=11, italic=True, color=LIGHT)
    # Right: 3 metric chips
    text_box(s, x=7.85, y=2.0, w=5.1, h=0.4,
             text="ХАРАКТЕРИСТИКИ",
             size=12, bold=True, color=TEAL)
    metrics = [
        ("текст → играбельный 3D-мир", MID),
        ("24 fps · реальное время", LIGHT),
        ("720p · несколько минут", GOLD),
    ]
    for i, (mtext, color) in enumerate(metrics):
        y = 2.5 + i * 0.85
        filled_rect(s, 7.85, y, 5.1, 0.7, color, radius=True, radius_adj=0.20)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=7.85, y=y, w=5.1, h=0.7,
                 text=mtext, size=18, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    # Anti-hype — shifted up + shortened to fit
    lesson_box(s, 7.85, 5.20, 5.1, 1.80,
               "Это не видеогенератор — это симулированное окружение. Промышленное использование: пока граничные случаи (прототипы игр, разведка локаций). Передний край — впереди.")
    speaker_notes(s, load_notes("s10"))


def build_s10a(p):
    """Russian context — локальное удобство vs frontier."""
    s = blank(p)
    text_box(s, x=0.55, y=0.30, w=12.3, h=0.65,
             text="Российский контекст: локальное удобство vs передний край",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.0, w=12.3, h=0.4,
             text="Российский GenAI функционален для конечного пользователя + масс-маркета премиум; но не на переднем крае в видео и музыке. Структурно (capex/data), а не идеологически.",
             size=12, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: real Russian AI image-gen ecosystem screenshot (Шедеврум web + mobile interfaces
    # с реальными AI-сгенерированными работами: осьминог-с-арбузом, крокодил-робот, дом-в-облаках)
    ocean_box(s, 0.55, 1.7, 5.5, 3.15)
    text_box(s, x=0.65, y=1.78, w=5.3, h=0.35,
             text="ШЕДЕВРУМ + KANDINSKY · ОФИЦИАЛЬНЫЕ ИНТЕРФЕЙСЫ",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s10a-kandinsky-vs-kling.png"),
              x=0.62, y=2.15, w=5.40, h=2.45)
    text_box(s, x=0.65, y=4.62, w=5.3, h=0.30,
             text="Источник: appleinsider.ru сравнение Шедеврум vs Kandinsky · 2023",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right: 4-card RU landscape — Ocean palette only (no red/green)
    grid_x = 6.40
    grid_y = 1.7
    grid_w = 6.55
    grid_h = 3.0
    cell_w = (grid_w - 0.15) / 2
    cell_h = (grid_h - 0.15) / 2
    areas = [
        ("Изображения", "Kandinsky 6.0 (MoE, бесплатно через GigaChat)\nYandex Шедеврум · YandexART 2.7",
         "конкурентоспособно", TEAL),
        ("Видео", "Kandinsky 5.0 (Apache 2.0)\nОтставание от Sora/Veo/Kling",
         "структурный разрыв", GOLD),
        ("Аудио", "SymFormer · SaluteSpeech\nYandex SpeechKit",
         "ниже ElevenLabs", MID),
        ("Правовое поле", "Минцифры 18.03.2026\nTDM* · маркировка · 01.09.2027",
         "в процессе", LIGHT),
    ]
    for i, (atitle, abody, atag, atag_color) in enumerate(areas):
        row = i // 2; col = i % 2
        x = grid_x + col * (cell_w + 0.15)
        y = grid_y + row * (cell_h + 0.15)
        ocean_box(s, x, y, cell_w, cell_h)
        text_box(s, x=x + 0.12, y=y + 0.08, w=cell_w - 0.24, h=0.30,
                 text=atitle, size=12, bold=True, color=DEEP)
        text_box(s, x=x + 0.12, y=y + 0.40, w=cell_w - 0.24, h=0.65,
                 text=abody, size=10, color=SLATE, line_spacing=1.30)
        chip_text_c = DEEP if atag_color == GOLD else WHITE
        chip(s, x + 0.12, y + cell_h - 0.45, cell_w - 0.24, 0.35, atag,
             fill=atag_color, color=chip_text_c, size=10, bold=True)
    # Inline glossary footer (TDM)
    text_box(s, x=6.40, y=4.78, w=6.55, h=0.30,
             text="* TDM = Text & Data Mining (закон об исключении для исследований)",
             size=9, italic=True, color=LIGHT, line_spacing=1.0)
    # Lesson box
    lesson_box(s, 0.55, 5.10, 12.3, 1.85,
               "Локальное удобство (бесплатно · RU-промпты · без VPN · рубли · правовой контур) ≠ frontier-уровень в видео и музыке. Концентрация R&D в США/Китае — структурное (capex GPU, видео-датасеты), не идеологическое.")
    speaker_notes(s, load_notes("s10a"))


def build_s11(p):
    """Personalisation at scale + Adobe Firefly."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Персонализация в масштабе + промышленное применение в Голливуде",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Adobe Firefly: $400 млн прямой выручки 2024-25. Lionsgate × Runway — усиление, а не замещение.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: Adobe Firefly коллаж логотипов + Lionsgate quote
    ocean_box(s, 0.55, 2.0, 7.0, 4.5)
    text_box(s, x=0.70, y=2.15, w=6.7, h=0.4,
             text="ADOBE FIREFLY · КОРПОРАТИВНЫЕ КЛИЕНТЫ",
             size=11, bold=True, color=TEAL)
    # Logos as text chips
    логотипы = ["Deloitte", "Tapestry", "Paramount+", "Pepsi", "dentsu", "Stagwell"]
    for i, logo in enumerate(логотипы):
        row = i // 3; col = i % 3
        x = 0.70 + col * 2.20
        y = 2.65 + row * 0.65
        filled_rect(s, x, y, 2.05, 0.5, WHITE, stroke=LIGHT, stroke_pt=1.0,
                    radius=True, radius_adj=0.15)
        text_box(s, x=x, y=y, w=2.05, h=0.5,
                 text=logo, size=14, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Lionsgate × Runway quote with официальное объявление о партнёрстве image
    filled_rect(s, 0.70, 4.10, 6.7, 2.2, GOLD_TINT, stroke=GOLD, stroke_pt=1.5,
                radius=True, radius_adj=0.05)
    # Real Runway × Lionsgate баннер с объявлением on left
    add_image(s, str(ASSETS / "screenshots/s11-lionsgate-runway.png"),
              x=0.85, y=4.20, w=2.10, h=1.40)
    text_box(s, x=3.10, y=4.20, w=4.20, h=0.4,
             text="LIONSGATE × RUNWAY · 18 СЕНТ 2024",
             size=10, bold=True, color=GOLD)
    text_box(s, x=3.10, y=4.65, w=4.20, h=1.4,
             text="«...экономия миллионов и миллионов долларов на препродакшене и постпродакшене»\n— Майкл Бёрнс, вице-председатель Lionsgate",
             size=11, italic=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=0.85, y=5.70, w=6.4, h=0.30,
             text="Источник: investors.lionsgate.com · звонок по итогам квартала, ноя 2024",
             size=9, italic=True, color=SLATE)
    # Right: 3 metric chips
    text_box(s, x=7.85, y=2.0, w=5.1, h=0.4,
             text="МЕТРИКИ 2026",
             size=11, bold=True, color=TEAL)
    metrics = [
        ("22 млрд+", "ассетов <2 лет"),
        ("$400 млн", "прямая выручка 2024-25"),
        ("40%", "видеорекламы — AI-сгенерирована (IAB 2026)"),
    ]
    for i, (val, label) in enumerate(metrics):
        y = 2.55 + i * 1.20
        filled_rect(s, 7.85, y, 5.1, 1.10, MID, radius=True, radius_adj=0.10)
        text_box(s, x=7.85, y=y + 0.10, w=5.1, h=0.55,
                 text=val, size=32, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=7.85, y=y + 0.70, w=5.1, h=0.35,
                 text=label, size=11, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Мини-провал
    lesson_box(s, 0.55, 6.6, 12.3, 0.85,
               "МИНИ-ПРОВАЛ: 86% закупщиков используют AI, но внедрение ≠ успех. Toys R Us Cannes 2024 — разворот тональности −10 п.п.")
    speaker_notes(s, load_notes("s11"))


def build_s13(p):
    build_section_divider(p, here_idx=2, title="AI ИЗМЕНИЛ",
                          frame_phrase="Обвал стоимости 100×-10 000× · скорость дни → секунды · новые vs потерянные роли.",
                          notes_slide_id="s13")


def build_s14(p):
    """Обвал стоимости 100×-10 000×."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Обвал стоимости: 100×–10 000× по типам ассетов",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Маржинальная стоимость генерации — на 2-4 порядка дешевле. Но коммерчески-безопасный корпоративный сегмент — отдельный.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Cost table — 4 rows
    table_x = 0.55
    table_y = 2.0
    table_w = 8.5
    row_h = 0.90
    # Header
    filled_rect(s, table_x, table_y, table_w, 0.5, DEEP)
    headers = [("АССЕТ", 0, 2.5), ("ДО", 2.5, 2.0), ("ПОСЛЕ", 4.5, 2.0), ("×", 6.5, 2.0)]
    for h, x_off, w in headers:
        text_box(s, x=table_x + x_off + 0.10, y=table_y + 0.10, w=w - 0.20, h=0.30,
                 text=h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("1 картинка", "$50-200 фриланс\n$25-100 стоковое фото", "$0-0.25", "200×-10 000×"),
        ("50 продуктовых картинок", "$1-25 тыс. фотосет", "$0-1.50", ">1 000×"),
        ("1 мин видео 720p", "$1-50 тыс. съёмка+постпродакшн", "~$6 Sora 2", "150×-8 000×"),
        ("Дубляж $/мин/язык", "$50-500 актёр+студия", "<$1 ElevenLabs", "50×-500×"),
    ]
    for i, (asset, before, after, mult) in enumerate(rows):
        y = table_y + 0.5 + i * row_h
        bg = SURFACE if i % 2 == 0 else WHITE
        filled_rect(s, table_x, y, table_w, row_h, bg,
                    stroke=LIGHT, stroke_pt=0.5)
        text_box(s, x=table_x + 0.10, y=y + 0.10, w=2.3, h=row_h - 0.20,
                 text=asset, size=12, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        text_box(s, x=table_x + 2.60, y=y + 0.10, w=1.8, h=row_h - 0.20,
                 text=before, size=11, italic=True, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        text_box(s, x=table_x + 4.60, y=y + 0.10, w=1.8, h=row_h - 0.20,
                 text=after, size=12, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        text_box(s, x=table_x + 6.60, y=y + 0.10, w=1.8, h=row_h - 0.20,
                 text=mult, size=13, bold=True, color=GOLD,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Right: $400M callout
    ocean_box(s, 9.25, 2.0, 3.65, 4.0, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.8)
    text_box(s, x=9.40, y=2.15, w=3.4, h=0.4,
             text="СРЕДНИЙ СЕГМЕНТ ≠ БЕСПЛАТНО",
             size=11, bold=True, color=GOLD)
    text_box(s, x=9.40, y=2.55, w=3.4, h=1.1,
             text="$400 млн",
             size=46, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=9.40, y=3.70, w=3.4, h=0.5,
             text="Adobe Firefly\nпрямая выручка фискальный год 2024-25",
             size=11, color=DEEP, bold=True,
             align=PP_ALIGN.CENTER, line_spacing=1.30)
    text_box(s, x=9.40, y=4.50, w=3.4, h=1.4,
             text="лицензированный корпус + процесс + интеграция = корпоративный SaaS-стек\n\nНижний сегмент вымывается, средний сегмент растёт.",
             size=10, italic=True, color=DEEP, line_spacing=1.30)
    # Footer
    text_box(s, x=0.55, y=6.75, w=12.3, h=0.4,
             text="Источник: ZSky AI, Sora 2 API Pricing, ElevenLabs pricing.",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    """Speed-collapse: дни → секунды."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Скорость: дни → секунды",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Итерационный цикл становится 10-100× плотнее → новые навыки человека.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 4 rows timer table
    rows = [
        ("Концепт-арт черновик", "дни (дизайнер-фрилансер)", "5-60 сек", "Midjourney / Flux / Imagen"),
        ("B-roll кадр", "часы съёмки + постпродакшн", "5-60 сек", "Veo / Sora"),
        ("Дубляж длинного формата", "недели в студии", "минуты", "ElevenLabs Dubbing Studio"),
        ("Исследование концепций", "полу-неделя · 3-5 вариантов", "минуты · 10×+", "итерация плотнее"),
    ]
    ocean_box(s, 0.55, 2.0, 12.3, 4.0)
    row_h = 0.85
    for i, (task, before, after, tool) in enumerate(rows):
        y = 2.20 + i * row_h
        # Task label
        text_box(s, x=0.80, y=y, w=2.8, h=row_h - 0.10,
                 text=task, size=14, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        # Before (Light teal — baseline, no red)
        filled_rect(s, 3.70, y + 0.10, 2.6, 0.55, SURFACE,
                    stroke=LIGHT, stroke_pt=1.2, radius=True, radius_adj=0.20)
        text_box(s, x=3.70, y=y + 0.10, w=2.6, h=0.55,
                 text=before, size=11, italic=True, color=LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Arrow
        right_arrow(s, 6.45, y + 0.20, 0.6, 0.35, fill=MID)
        # After (gold accent — speed win)
        filled_rect(s, 7.20, y + 0.10, 2.0, 0.55, GOLD,
                    radius=True, radius_adj=0.20)
        text_box(s, x=7.20, y=y + 0.10, w=2.0, h=0.55,
                 text=after, size=15, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Tool
        text_box(s, x=9.35, y=y + 0.10, w=3.5, h=0.55,
                 text=tool, size=11, italic=True, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Anchor
    lesson_box(s, 0.55, 6.2, 12.3, 1.05,
               "Инженерный урок: итерационный цикл 10-100× плотнее → новые навыки человека — формулировать «годен/не годен» быстрее, потому что новых вариантов появляется больше.")
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    """New professions."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Новые профессии: prompt-инженер · AI-режиссёр · специалист по AI-процессам",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Специализированный класс между AI-инструментом и итоговым продуктом (client deliverable).",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: Upwork screenshot мокап
    ocean_box(s, 0.55, 2.0, 5.5, 4.0)
    text_box(s, x=0.75, y=2.15, w=5.1, h=0.4,
             text="UPWORK · КАТЕГОРИЯ AI/ML",
             size=11, bold=True, color=TEAL)
    filled_rect(s, 0.75, 2.65, 5.1, 0.5, MID, radius=True, radius_adj=0.10)
    text_box(s, x=0.85, y=2.65, w=4.9, h=0.5,
             text="Поиск: AI / ML / prompt-инжиниринг",
             size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Mock results list
    jobs = [
        "Prompt-инженер · Midjourney эксперт · $35-65/час",
        "AI-режиссёр / супервайзер · $50-90/час",
        "Специалист по AI-процессам · корп. · $60-120/час",
        "Супервайзер континьюити · видео · $40-75/час",
    ]
    for i, job in enumerate(jobs):
        y = 3.30 + i * 0.55
        filled_rect(s, 0.75, y, 5.1, 0.45, SURFACE,
                    stroke=LIGHT, stroke_pt=0.5, radius=True, radius_adj=0.18)
        text_box(s, x=0.85, y=y, w=4.9, h=0.45,
                 text=job, size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    # Metrics chip at bottom of card
    chip(s, 0.75, 5.55, 5.1, 0.40, "Upwork +70% год к году · 52% объёма услуг — AI",
         fill=GOLD_TINT, stroke=GOLD, color=DEEP, size=10, bold=True)
    # Right: 4 role-cards
    text_box(s, x=6.35, y=2.0, w=6.6, h=0.4,
             text="4 НОВЫЕ РОЛИ",
             size=11, bold=True, color=TEAL)
    roles = [
        ("Prompt-инженер / AI-художник",
         "Формирует промпты + постобработка → готовый к промышленному применению результат",
         "$25-80/час"),
        ("AI-режиссёр / AI-музыкальный продюсер",
         "Супервайзер вывода моделей: итерации + постобработка + мультимодальность",
         "Аналог арт-директора"),
        ("Специалист по AI-процессам",
         "Интегратор AI-инструментов в производственные процессы студий",
         "развёртывания Adobe Firefly Foundry"),
        ("Супервайзер непрерывности",
         "Проверка непрерывности персонажа/сцены в мульти-кадровых последовательностях",
         "Контроль дрейфа между сценами"),
    ]
    for i, (rtitle, rbody, rextra) in enumerate(roles):
        y = 2.5 + i * 0.95
        ocean_box(s, 6.35, y, 6.6, 0.85)
        text_box(s, x=6.50, y=y + 0.05, w=6.3, h=0.30,
                 text=rtitle, size=12, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=6.50, y=y + 0.34, w=6.3, h=0.28,
                 text=rbody, size=10, color=SLATE, line_spacing=1.15)
        text_box(s, x=6.50, y=y + 0.60, w=6.3, h=0.22,
                 text=rextra, size=10, italic=True, color=GOLD, bold=True,
                 line_spacing=1.0)
    # Lesson box — repositioned to fit fully
    lesson_box(s, 0.55, 6.30, 12.3, 1.00,
               "Новые роли — между AI-инструментом и итоговым продуктом для клиента. Растут быстро, но меньше замещённого класса.")
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """Замещение."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Замещение: −17% графический дизайн · Shutterstock единичные значения",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Снижение зарплат снизу — структурная, не временная. Консолидация индустрии как ответ.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: big -17% callout (3 columns shrunk to make room for Урок full-width below)
    ocean_box(s, 0.55, 2.0, 4.0, 4.0)
    text_box(s, x=0.70, y=2.10, w=3.7, h=0.30,
             text="UPWORK · ГРАФИЧЕСКИЙ ДИЗАЙН · ГОД К ГОДУ",
             size=10, bold=True, color=TEAL, line_spacing=1.0)
    text_box(s, x=0.70, y=2.50, w=3.7, h=1.30,
             text="−17%",
             size=68, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.70, y=3.85, w=3.7, h=0.30,
             text="работ в графическом дизайне",
             size=11, italic=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, x=0.70, y=4.25, w=3.7, h=1.65,
             text="40% работ копирайтеров $10-19/час — распознан как AI · <10% в $60+/час · снижение зарплат снизу",
             size=10, color=SLATE, line_spacing=1.30, align=PP_ALIGN.CENTER)
    # Middle: Shutterstock licensing timeline
    ocean_box(s, 4.75, 2.0, 4.0, 4.0)
    text_box(s, x=4.90, y=2.10, w=3.7, h=0.30,
             text="SHUTTERSTOCK · ЛИЦЕНЗИРОВАНИЕ ДЛЯ AI",
             size=10, bold=True, color=TEAL)
    # Timeline bars
    years_data = [("2023", 104, MID), ("2024", 138, LIGHT), ("2027", 250, GOLD)]
    max_val = 250
    for i, (year, val, color) in enumerate(years_data):
        y = 2.55 + i * 0.85
        bar_w = (val / max_val) * 3.5
        text_box(s, x=4.90, y=y, w=0.6, h=0.35,
                 text=year, size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        filled_rect(s, 5.55, y, bar_w, 0.45, color, radius=True, radius_adj=0.20)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=5.65, y=y, w=bar_w - 0.1, h=0.45,
                 text=f"${val} млн", size=13, bold=True, color=text_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=4.90, y=5.45, w=3.7, h=0.30,
             text="разворот: фото → обучающие данные для AI",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right: Getty merger + SAG-AFTRA
    ocean_box(s, 8.95, 2.0, 4.0, 4.0)
    text_box(s, x=9.10, y=2.10, w=3.7, h=0.30,
             text="ОТВЕТ ИНДУСТРИИ",
             size=10, bold=True, color=TEAL)
    # Getty/Shutterstock merger
    filled_rect(s, 9.10, 2.50, 3.7, 0.90, MID, radius=True, radius_adj=0.10)
    text_box(s, x=9.20, y=2.55, w=3.5, h=0.35,
             text="Слияние Getty + Shutterstock", size=10, bold=True, color=WHITE)
    text_box(s, x=9.20, y=2.92, w=3.5, h=0.45,
             text="$3.7B · Jan 2025 · Оборонительная консолидация",
             size=9, italic=True, color=WHITE, line_spacing=1.20)
    # SAG-AFTRA chip (Screen Actors Guild) + WGA (Writers Guild)
    filled_rect(s, 9.10, 3.50, 3.7, 0.90, LIGHT, radius=True, radius_adj=0.10)
    text_box(s, x=9.20, y=3.55, w=3.5, h=0.35,
             text="SAG-AFTRA + WGA*", size=10, bold=True, color=WHITE)
    text_box(s, x=9.20, y=3.92, w=3.5, h=0.45,
             text="Оговорка о цифровых репликах · 2026 · продление на 4 года",
             size=9, italic=True, color=WHITE, line_spacing=1.20)
    # Voice actors
    filled_rect(s, 9.10, 4.50, 3.7, 0.90, GOLD, radius=True, radius_adj=0.10)
    text_box(s, x=9.20, y=4.55, w=3.5, h=0.35,
             text="Актёры озвучки · в мире", size=10, bold=True, color=DEEP)
    text_box(s, x=9.20, y=4.92, w=3.5, h=0.45,
             text="Коммодити-дубляж вытеснен · ElevenLabs в корпоративном секторе",
             size=9, italic=True, color=DEEP, line_spacing=1.20)
    # Glossary footer (smaller, bottom of right column)
    text_box(s, x=9.10, y=5.55, w=3.7, h=0.30,
             text="* SAG-AFTRA = Гильдия киноактёров · WGA = Гильдия сценаристов",
             size=8, italic=True, color=LIGHT, align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Lesson — full width below all 3 columns
    lesson_box(s, 0.55, 6.20, 12.3, 1.10,
               "Структурно: AI вымывает нижний сегмент, оставляет верхний сегмент защищённым. Clauses не покрывают bottom фриланс.")
    speaker_notes(s, load_notes("s17"))


def build_s19(p):
    build_section_divider(p, here_idx=3, title="AI СЛОМАЛ",
                          frame_phrase="12 кейсов: авторское право × 4 + дипфейки × 2 + slop · фейковые авторы · негативная реакция · замещение + таксономия.",
                          notes_slide_id="s19")


def build_s20(p):
    """Copyright 4 categories таксономия."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="«AI и авторское право» — 4 разных категории исков",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Не один юридический вопрос, а четыре. С разной правовой логикой и исходами.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 2x2 matrix
    grid_x = 0.55
    grid_y = 2.0
    grid_w = 8.5
    grid_h = 4.3
    cell_w = (grid_w - 0.20) / 2
    cell_h = (grid_h - 0.20) / 2
    cats = [
        ("1. Парсинг для обучения", "СО СТОРОНЫ ВХОДА",
         "AI-компания собрала корпус БЕЗ лицензии",
         "→ NYT против OpenAI · Andersen", MID),
        ("2. Сходство вывода", "СО СТОРОНЫ ВЫХОДА",
         "Модель воспроизводит охраняемый контент дословно",
         "→ дословное цитирование NYT · DMCA", LIGHT),
        ("3. Подражание стилю", "«В СТИЛЕ КОНКРЕТНОГО ХУДОЖНИКА»",
         "Коллективный иск художников · стиль не охраняется авторским правом, но DMCA + публичные права",
         "→ Andersen против Stability/MJ/Deviant", GOLD),
        ("4. Голос/образ", "ПРАВО НА ОБРАЗ",
         "Использование голоса/образа без согласия",
         "→ ScarJo против OpenAI · SAG-AFTRA · Корея", TEAL),
    ]
    for i, (title, subtitle, body, cases, accent) in enumerate(cats):
        row = i // 2; col = i % 2
        x = grid_x + col * (cell_w + 0.20)
        y = grid_y + row * (cell_h + 0.20)
        ocean_box(s, x, y, cell_w, cell_h)
        # Accent corner
        filled_rect(s, x, y, cell_w, 0.55, accent, radius=True, radius_adj=0.10)
        accent_text = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=y + 0.08, w=cell_w - 0.30, h=0.40,
                 text=title, size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.18, y=y + 0.65, w=cell_w - 0.36, h=0.35,
                 text=subtitle, size=10, bold=True, italic=True, color=TEAL)
        text_box(s, x=x + 0.18, y=y + 1.00, w=cell_w - 0.36, h=0.7,
                 text=body, size=12, color=DEEP, line_spacing=1.30)
        text_box(s, x=x + 0.18, y=y + 1.65, w=cell_w - 0.36, h=0.4,
                 text=cases, size=10, italic=True, color=GOLD, bold=True,
                 line_spacing=1.20)
    # Right: lesson box
    lesson_box(s, 9.30, 2.0, 3.65, 4.3,
               "«Авторское право AI» — НЕ один вопрос, а 4 разных категории риска.\n\nСмотри, какая из 4 применима к твоему процесс.")
    # Footer
    text_box(s, x=0.55, y=6.55, w=12.3, h=0.35,
             text="Каждая из 4 категорий раскрыта далее на отдельном эталонном кейсе.",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    speaker_notes(s, load_notes("s20"))


def case_slide_template(p, slide_id, title, assertion_body, screenshot_path,
                        timeline_events, lesson_text, screenshot_label=None,
                        emphasis_block=None, glossary=None):
    """Generic case-slide builder for s21-s25, s27. Phase 6+7 v2.

    screenshot_path — abs path to real news-card PNG (REQUIRED).
    emphasis_block — optional dict for slide-specific emphasis layout:
        {"kind": "big_number", "value": "20M", "label": "ChatGPT logs"}
        {"kind": "verdict_badge", "value": "STABILITY ВЫИГРАЛ", "caption": "..."}
        {"kind": "trial_chip", "value": "8 СЕНТ 2026", "label": "ДАТА СУДА"}
    glossary — short инлайн-пояснение (e.g. «MTD = motion to dismiss»)
    """
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text=title, size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text=assertion_body, size=13, italic=True, color=MID, align=PP_ALIGN.LEFT,
             line_spacing=1.30)
    # Left: real news-screenshot embedded
    ocean_box(s, 0.55, 2.0, 6.0, 4.0)
    if screenshot_label:
        text_box(s, x=0.70, y=2.05, w=5.7, h=0.3,
                 text=screenshot_label, size=10, bold=True, color=TEAL)
    if screenshot_path and Path(screenshot_path).exists():
        add_image(s, screenshot_path, x=0.70, y=2.40, w=5.70, h=3.50)
    # Right: timeline events
    ocean_box(s, 6.85, 2.0, 6.10, 4.0)
    text_box(s, x=7.00, y=2.10, w=5.8, h=0.4,
             text="ХРОНОЛОГИЯ", size=11, bold=True, color=TEAL)
    for i, (date, event, color) in enumerate(timeline_events):
        y = 2.60 + i * 0.78
        # Date chip
        filled_rect(s, 7.00, y, 1.5, 0.5, color, radius=True, radius_adj=0.20)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=7.00, y=y, w=1.5, h=0.5, text=date,
                 size=10, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Event text
        text_box(s, x=8.60, y=y, w=4.3, h=0.5, text=event,
                 size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    # Emphasis block (slide-specific diversification)
    if emphasis_block:
        eb = emphasis_block
        if eb.get("kind") == "big_number":
            # Big number callout below timeline
            filled_rect(s, 6.85, 5.65, 6.10, 0.45, GOLD,
                        radius=True, radius_adj=0.20)
            text_box(s, x=6.85, y=5.65, w=6.10, h=0.45,
                     text=f"{eb['value']}  ·  {eb['label']}",
                     size=12, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "verdict_badge":
            filled_rect(s, 6.85, 5.65, 6.10, 0.45, TEAL,
                        radius=True, radius_adj=0.20)
            text_box(s, x=6.85, y=5.65, w=6.10, h=0.45,
                     text=f"{eb['value']}",
                     size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "trial_chip":
            filled_rect(s, 6.85, 5.65, 6.10, 0.45, MID,
                        radius=True, radius_adj=0.20)
            text_box(s, x=6.85, y=5.65, w=6.10, h=0.45,
                     text=f"{eb['label']} · {eb['value']}",
                     size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "settlement_matrix":
            # 3-major × 2-defendant matrix (2 rows of 3, или 1 row of 3 для legacy)
            label_y = 5.40
            filled_rect(s, 6.85, label_y, 6.10, 0.28, SURFACE,
                        radius=True, radius_adj=0.30)
            text_box(s, x=6.85, y=label_y, w=6.10, h=0.28,
                     text=eb.get("title", "КРУПНЫЕ ЛЕЙБЛЫ · СТАТУС УРЕГУЛИРОВАНИЯ"),
                     size=9, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cells = eb.get("cells", [])
            # Layout: 3 columns x N rows
            ncols = 3
            nrows = (len(cells) + ncols - 1) // ncols
            cw = 6.10 / ncols
            ch = 0.32
            for ci, (label, color) in enumerate(cells):
                row = ci // ncols
                col = ci % ncols
                cx = 6.85 + col * cw
                cy = label_y + 0.32 + row * (ch + 0.05)
                filled_rect(s, cx + 0.05, cy, cw - 0.10, ch,
                            color, radius=True, radius_adj=0.30)
                ct = WHITE if color != GOLD else DEEP
                text_box(s, x=cx, y=cy, w=cw, h=ch,
                         text=label, size=8, bold=True, color=ct,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif eb.get("kind") == "fair_use_factors":
            label_y = 5.50
            filled_rect(s, 6.85, label_y, 6.10, 0.30, SURFACE,
                        radius=True, radius_adj=0.30)
            text_box(s, x=6.85, y=label_y, w=6.10, h=0.30,
                     text="«ДОБРОСОВЕСТНОЕ ИСПОЛЬЗОВАНИЕ» 4-ФАКТОРНЫЙ ТЕСТ (Warhol против Goldsmith)",
                     size=9, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            factors = eb.get("factors", [])
            fw = 6.10 / max(1, len(factors))
            for fi, (label, status_color) in enumerate(factors):
                fx = 6.85 + fi * fw
                filled_rect(s, fx + 0.05, label_y + 0.35, fw - 0.10, 0.30,
                            status_color, radius=True, radius_adj=0.30)
                ct = WHITE if status_color != GOLD else DEEP
                text_box(s, x=fx, y=label_y + 0.35, w=fw, h=0.30,
                         text=label, size=9, bold=True, color=ct,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Optional glossary footnote
    if glossary:
        text_box(s, x=6.85, y=6.05, w=6.10, h=0.25,
                 text=glossary, size=8, italic=True, color=LIGHT,
                 align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Lesson box at bottom
    lesson_box(s, 0.55, 6.30, 12.3, 1.05, lesson_text)
    speaker_notes(s, load_notes(slide_id))


def build_s21(p):
    # Emphasis: 20M ChatGPT logs (big number callout)
    case_slide_template(p, "s21",
        title="NYT против OpenAI — обучение + сходство вывода (Кейс 1)",
        assertion_body="NYT против OpenAI (Дек 2023): суд обязал OpenAI выдать 20 млн ChatGPT-логов. Дедлайн упрощённого решения суда (SJ*) — 2 апр 2026. Теория дословного цитирования.",
        screenshot_path=str(ASSETS / "screenshots/s21-nyt-bloomberg.png"),
        screenshot_label="THE NEW YORK TIMES · 27 ДЕК 2023",
        timeline_events=[
            ("Дек 2023", "NYT подал иск", MID),
            ("2024-25", "Истребование доказательств + борьба за процессуальные действия", LIGHT),
            ("2 апр 2026", "Дедлайн упрощённого решения суда (SJ)", GOLD),
        ],
        emphasis_block={"kind": "big_number", "value": "20 000 000", "label": "ChatGPT-логов под истребование доказательств"},
        glossary="* SJ = summary judgment (упрощённое решение суда без полного процесса)",
        lesson_text="Если модель может процитировать твой обучающий корпус дословно — это НЕ «добросовестное использование», это доказательство нарушения. Проверка сходства результата с обучающими данными — обязательна."
    )


def build_s22(p):
    # Emphasis: UK verdict badge "STABILITY ВЫИГРАЛ"
    case_slide_template(p, "s22",
        title="Getty против Stability — победа в UK vs ожидание в US (Кейс 2)",
        assertion_body="Высокий суд UK 04.11.2025 — Stability выиграл основные требования (weights ≠ копия по CDPA*). Кейс в US — motion to dismiss** 10.02.2026.",
        screenshot_path=str(ASSETS / "screenshots/s22-getty-bird.png"),
        screenshot_label="THE VERGE · GETTY vs STABILITY · ФЕВ 2023",
        timeline_events=[
            ("04 ноя 2025", "UK · Stability выиграл основные требования", TEAL),
            ("10 фев 2026", "Дело в США · motion to dismiss ожидается", GOLD),
            ("Ожидается", "Товарный знак + passing-off — отдельные требования", LIGHT),
        ],
        emphasis_block={"kind": "verdict_badge",
                        "value": "UK: STABILITY ВЫИГРАЛ основные требования по CDPA"},
        glossary="* CDPA = UK Copyright, Designs and Patents Act 1988  ·  ** MTD = motion to dismiss",
        lesson_text="Юрисдикции расходятся — то, что законно в UK по CDPA, НЕЗАКОННО в US по «добросовестному использованию». Для глобального развёртывания проверяй обе."
    )


def build_s23(p):
    # Emphasis: trial date prominent
    case_slide_template(p, "s23",
        title="Andersen против Stability/MJ/Deviant — подражание стилю (Кейс 3)",
        assertion_body="Коллективный иск художников. Motion to dismiss* отклонён в авг 2024 → истребование доказательств. Третье изменённое заявление (amended complaint) — 27 фев 2026. Суд — 8 сент 2026.",
        screenshot_path=str(ASSETS / "screenshots/s23-andersen-docket.png"),
        screenshot_label="KELLY McKERNAN (ИСТЕЦ) · WIKIMEDIA COMMONS",
        timeline_events=[
            ("Янв 2023", "Коллективный иск подан (10 художников)", MID),
            ("Авг 2024", "motion to dismiss отклонён (судья Orrick) → истребование доказательств", LIGHT),
            ("27 фев 2026", "3-е изменённое заявление", TEAL),
        ],
        emphasis_block={"kind": "trial_chip", "label": "ДАТА СУДА", "value": "8 СЕНТ 2026"},
        glossary="* MTD = motion to dismiss",
        lesson_text="Подражание стилю «в стиле [конкретного художника]» — НЕБЕЗОПАСНО, даже если стиль не охраняется авторским правом. Коллективные иски проходят motion to dismiss на DMCA + публичные права."
    )


def build_s24(p):
    # Emphasis: 3 major × 2 defendant matrix
    case_slide_template(p, "s24",
        title="RIAA против Suno/Udio — лицензирование под давлением исков (Кейс 4)",
        assertion_body="RIAA подала иск 24.06.2024. UMG урегулировала с Udio 29.10.2025; Warner урегулировала с Suno сент 2025. Sony — активно судится с обоими.",
        screenshot_path=str(ASSETS / "screenshots/s24-riaa-suno.png"),
        screenshot_label="BILLBOARD · MAJOR LABEL LAWSUIT · 24 ИЮНЯ 2024",
        timeline_events=[
            ("24 июн 2024", "RIAA подаёт иск против Suno + Udio", MID),
            ("29 окт 2025", "UMG × Udio урегулирование → совместная платформа", TEAL),
            ("сент 2025", "Warner × Suno урегулирование (отчисления + доля)", TEAL),
            ("июл 2026", "Suno SJ — Sony активно судится с обоими", GOLD),
        ],
        emphasis_block={"kind": "settlement_matrix",
                        "title": "3 КРУПНЫХ ЛЕЙБЛА × 2 ОТВЕТЧИКА — статус",
                        "cells": [
                            ("UMG × Udio: урегулировано", TEAL),
                            ("UMG × Suno: переговоры", LIGHT),
                            ("Warner × Suno: урегулировано", TEAL),
                            ("Warner × Udio: тяжба", GOLD),
                            ("Sony × Suno: тяжба", GOLD),
                            ("Sony × Udio: тяжба", GOLD),
                        ]},
        glossary="* UMG = Universal Music Group (один из 3 крупных лейблов «Большой тройки»)",
        lesson_text="Лицензирование под давлением исков — фактический исход: 4 из 6 комбинаций иск-ответчик урегулированы или в переговорах. Это новый слой бизнес-модели, а не «запрет всей AI-музыки»."
    )


def build_s25(p):
    # Emphasis: «добросовестное использование» — 4-факторный тест breakdown
    case_slide_template(p, "s25",
        title="Thomson Reuters против Ross — первый отказ США в «добросовестном использовании» (Кейс 5)",
        assertion_body="Фев 2025, судья Bibas: 2200/3000 headnotes — нарушение, «добросовестное использование» отклонено. Оговорка: Ross — не-генеративный AI.",
        screenshot_path=str(ASSETS / "screenshots/s25-thomson-reedsmith.png"),
        screenshot_label="DAVIS WRIGHT TREMAINE · ROSS RULING · ФЕВ 2025",
        timeline_events=[
            ("Фев 2025", "Bibas · первый отказ США в «добросовестном использовании»", GOLD),
            ("2200/3000", "headnotes — нарушение", LIGHT),
            ("Ожидается", "Тестовые кейсы для LLM/diffusion впереди", MID),
        ],
        emphasis_block={"kind": "fair_use_factors",
                        "factors": [
                            ("1 · цель", GOLD),
                            ("2 · характер", GOLD),
                            ("3 · объём", GOLD),
                            ("4 · рынок", GOLD),
                        ]},
        glossary="Caveat: Ross — non-generative AI (legal search). Generative LLM/diffusion test cases ещё впереди.",
        lesson_text="«Добросовестное использование» — НЕ дефолт. Тестовые кейсы для LLM/diffusion впереди. НЕ строй продуктовую дорожную карту на предположении, что «добросовестное использование» сработает как защита."
    )


def build_s26(p):
    """Arup deepfake — special schema (attack diagram). NOTE: lesson_box already has УРОК prefix so we strip it from text."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Arup CFO-дипфейк — мошенничество на $25.6 млн (Кейс 6)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Гонконг, янв 2024 · Финансист на видеозвонке с дипфейк-CFO + коллегами → 15 транзакций.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: CNN главное изображение — actual published photo (hands on laptop / hacker imagery)
    ocean_box(s, 0.55, 2.0, 4.5, 4.0)
    text_box(s, x=0.70, y=2.08, w=4.2, h=0.30,
             text="CNN · МОШЕННИЧЕСТВО НА $25 МЛН В ГОНКОНГЕ · 16 МАЯ 2024", size=10, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s26-arup-cnn.png"),
              x=0.70, y=2.40, w=4.20, h=3.30)
    text_box(s, x=0.70, y=5.75, w=4.2, h=0.30,
             text="Arup · британская инженерная фирма · Sydney Opera House",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Right: attack diagram — 5 stages (Gold replaces Red on final step)
    text_box(s, x=5.35, y=2.0, w=7.6, h=0.4,
             text="СЦЕНАРИЙ АТАКИ", size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    stages = [
        ("1", "Письмо от «CFO»", MID),
        ("2", "Приглашение на видео", LIGHT),
        ("3", "Звонок с дипфейками", TEAL),
        ("4", "15 транзакций", MID),
        ("5", "$25.6 млн утрачено", GOLD),
    ]
    stage_y = 2.55
    stage_w = 1.30
    arrow_w = 0.20
    total_units = len(stages) * stage_w + (len(stages) - 1) * arrow_w
    start_x = 5.35 + (7.6 - total_units) / 2
    for i, (num, label, color) in enumerate(stages):
        x = start_x + i * (stage_w + arrow_w)
        filled_rect(s, x, stage_y, stage_w, 2.5, color, radius=True, radius_adj=0.10)
        text_color = WHITE if color != GOLD else DEEP
        text_box(s, x=x, y=stage_y + 0.20, w=stage_w, h=0.6, text=num,
                 size=32, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(s, x=x + 0.05, y=stage_y + 1.0, w=stage_w - 0.10, h=1.4,
                 text=label, size=10, color=text_color, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)
        if i < len(stages) - 1:
            right_arrow(s, x + stage_w + 0.02, stage_y + 1.0,
                        arrow_w - 0.04, 0.5, fill=DEEP)
    # Lesson
    lesson_box(s, 0.55, 6.4, 12.3, 0.95,
               "Видеозвонок ≠ подтверждение личности в 2024+. Финансовые транзакции требуют проверка через независимый канал — обратный звонок по известному номеру, многофакторная аутентификация, документированный процесс.")
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """Korea deepfake crisis — text-only, NO deepfake visuals."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Кризис со школьницами в Корее: дипфейки (кейс 7)",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Авг 2024 · >200 Telegram-чатов из селфи одноклассниц/учительниц · 4× к 2023.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: PBS NewsHour real photo — masked protestors with «반복되는 딥페이크 성범죄 국가도 공범이다»
    # banner (=«Repeated deepfake sex crimes, the state is an accomplice»). NO deepfake visuals.
    ocean_box(s, 0.55, 2.0, 5.0, 4.0)
    text_box(s, x=0.70, y=2.08, w=4.7, h=0.30,
             text="PBS NEWSHOUR · ПРОТЕСТ В КОРЕЕ · 2024", size=10, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s27-korea-npr.png"),
              x=0.70, y=2.40, w=4.70, h=3.50)
    # Right: 4 data cards
    text_box(s, x=5.85, y=2.0, w=7.1, h=0.4,
             text="ЦИФРЫ КРИЗИСА", size=11, bold=True, color=TEAL)
    stats = [
        (">200", "Telegram-чатов с дипфейк-порно"),
        ("6 500", "запросов на удаление, янв-июл 2024 (4× к 2023)"),
        ("74%", "подозреваемых — 10-19 лет"),
        ("793 / 16", "сообщений / уголовных дел (2021 — июл 2024)"),
    ]
    for i, (val, label) in enumerate(stats):
        y = 2.55 + i * 0.95
        ocean_box(s, 5.85, y, 7.1, 0.85)
        text_box(s, x=6.00, y=y + 0.10, w=2.0, h=0.65,
                 text=val, size=28, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=8.10, y=y + 0.10, w=4.7, h=0.65,
                 text=label, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.30)
    # Lesson
    lesson_box(s, 0.55, 6.4, 12.3, 0.95,
               "Доступная возможность + слабый правоприменительный контроль = массовый коллективный вред. Для AI-инструментов для конечных пользователей обязателен слой безопасности (детекция NSFW-контента + верификация возраста + процесс обработки жалоб) ДО запуска.")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """Slop + коллапс моделей."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Slop + коллапс моделей · Google AI Overviews (Кейс 8)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="«Put glue on pizza» + «eat one rock per day». Shumailov, Nature 2024: рекурсивное обучение → деградация.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: 2 Google AI Overview screenshots — Ocean palette
    ocean_box(s, 0.55, 2.0, 7.0, 4.2)
    text_box(s, x=0.70, y=2.10, w=6.7, h=0.4,
             text="GOOGLE AI OVERVIEW · реальные ответы (май 2024)",
             size=10, bold=True, color=TEAL)
    # Screenshot 1
    filled_rect(s, 0.70, 2.55, 3.30, 1.7, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.06)
    text_box(s, x=0.85, y=2.65, w=3.10, h=0.40,
             text="«how to keep cheese\nfrom sliding off pizza»",
             size=10, bold=True, color=DEEP, line_spacing=1.20)
    # Gold accent (anti-pattern flag) — not red
    filled_rect(s, 0.78, 3.18, 0.04, 0.95, fill=GOLD)
    text_box(s, x=0.95, y=3.20, w=3.00, h=1.0,
             text="→ AI Overview: «add ⅛ cup of non-toxic glue to the sauce»",
             size=10, italic=True, color=DEEP, bold=True, line_spacing=1.30)
    text_box(s, x=0.85, y=4.05, w=3.10, h=0.20,
             text="source: шутка из Reddit (11 years old)",
             size=9, italic=True, color=LIGHT)
    # Screenshot 2
    filled_rect(s, 4.20, 2.55, 3.20, 1.7, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.06)
    text_box(s, x=4.35, y=2.65, w=2.9, h=0.40,
             text="«how many rocks should\nI eat per day?»",
             size=10, bold=True, color=DEEP, line_spacing=1.20)
    filled_rect(s, 4.28, 3.18, 0.04, 0.95, fill=GOLD)
    text_box(s, x=4.45, y=3.20, w=2.80, h=1.0,
             text="→ AI Overview: «at least one small rock per day»",
             size=10, italic=True, color=DEEP, bold=True, line_spacing=1.30)
    text_box(s, x=4.35, y=4.05, w=2.9, h=0.20,
             text="source: The Onion (satire)",
             size=9, italic=True, color=LIGHT)
    # Below: Nature paper card with real Fig 1 (perplexity histograms коллапс моделей)
    filled_rect(s, 0.70, 4.50, 6.7, 1.55, MID, radius=True, radius_adj=0.08)
    # Real Nature paper Fig 1 — Shumailov 2024 коллапс моделей figure on left
    add_image(s, str(ASSETS / "screenshots/s28-ai-overview-glue.png"),
              x=0.80, y=4.60, w=1.40, h=1.35)
    text_box(s, x=2.30, y=4.60, w=5.0, h=0.40,
             text="NATURE · vol 631 · p 755-759 (2024)",
             size=11, bold=True, color=GOLD)
    text_box(s, x=2.30, y=5.05, w=5.0, h=0.55,
             text="Shumailov et al. — «AI-модели коллапсируют при обучении на рекурсивно-сгенерированных данных»",
             size=11, color=WHITE, bold=True, line_spacing=1.25)
    text_box(s, x=0.85, y=5.65, w=6.4, h=0.35,
             text="MAD: Model Autophagy Disorder (расстройство аутофагии моделей)",
             size=11, italic=True, color=WHITE, line_spacing=1.0)
    # Right: lesson
    lesson_box(s, 7.70, 2.0, 5.25, 4.2,
               "Качество источников важнее объёма.\n\nМодель на шутках Reddit без фильтрации проигрывает модели на курируемом датасете — даже если курируемый в 10× меньше.\n\nЭто объясняет, почему Adobe Firefly работает: лицензированный корпус, а не данные, спарсенные из веба.")
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """SI — фейковые авторы + Amazon."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Sports Illustrated фейковые авторы + Amazon фейковые книги (Кейс 9)",
             size=22, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="SI, ноя 2023 · статьи с AI-сгенерированными фото профилей. Authors Guild: всплеск фейковых книг на Amazon Kindle с AI-псевдонимами.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: CNN article screenshot — actual Drew Ortiz фейковый профиль from SI website
    # (Sports Illustrated bio page with AI face + AI-сгенерированный outdoors-reviewer text)
    ocean_box(s, 0.55, 2.0, 5.5, 4.0)
    text_box(s, x=0.70, y=2.08, w=5.2, h=0.30,
             text="CNN · SPORTS ILLUSTRATED ФЕЙКОВЫЕ АВТОРЫ · 27 НОЯ 2023", size=10, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s29-si-futurism.png"),
              x=0.70, y=2.40, w=5.20, h=2.65)
    chip(s, 0.70, 5.20, 5.2, 0.4,
         "«Drew Ortiz» — полностью вымышленный · AI-портрет + AI-текст",
         fill=GOLD, color=DEEP, size=10, bold=True)
    # Right: Amazon Kindle Authors Guild data (no fabricated 19/100 number)
    ocean_box(s, 6.30, 2.0, 6.65, 4.0)
    text_box(s, x=6.45, y=2.10, w=6.3, h=0.4,
             text="AMAZON KINDLE 2023-24 · ГИЛЬДИЯ АВТОРОВ",
             size=11, bold=True, color=TEAL)
    # Big icon/headline replaces fabricated 19/100 stat
    text_box(s, x=6.45, y=2.55, w=6.3, h=1.0,
             text="AI-ПСЕВДОНИМЫ",
             size=44, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=6.45, y=3.55, w=6.3, h=0.5,
             text="всплеск фейковых книг (Authors Guild 2023-24)",
             size=13, italic=True, color=DEEP, align=PP_ALIGN.CENTER, bold=True)
    text_box(s, x=6.45, y=4.10, w=6.3, h=1.85,
             text="AI-сгенерированные knockoffs выдают себя за реальных jazz-figures, финансовых консультантов\n\n«Frank Gioia» · «Ted Alkyer» — fakes реальных jazz-персоналий\n\nAmazon ввёл лимит 3 книги/день + AI-disclosure",
             size=11, color=SLATE, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.40)
    # Lesson
    lesson_box(s, 0.55, 6.20, 12.3, 1.0,
               "Накопленное доверие — ключевой актив бренда. AI-псевдонимы разрушают его моментально. Если публикуешь под именем — имя должно быть реальным человеком ИЛИ авторство AI должно быть явно раскрыто.")
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """Toys R Us / Coca-Cola разворот тональности — real ad screenshots + delta."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Маркетинговый провал — Toys R Us + Coca-Cola (Кейс 10)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Toys R Us Cannes Lions 2024 — разворот тональности −10 п.п. Coca-Cola Holidays 2024 AI-реклама — вирусный негатив.",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # Left: Toys R Us official Sora ad still
    ocean_box(s, 0.55, 2.0, 6.0, 2.55)
    text_box(s, x=0.70, y=2.10, w=5.7, h=0.4,
             text="TOYS R US STUDIOS · SORA-РЕКЛАМА · CANNES 2024",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s30-toysrus-cannes.png"),
              x=0.70, y=2.45, w=5.70, h=2.00)
    # Right: Coca-Cola Holidays Are Coming AI-реклама still
    ocean_box(s, 6.85, 2.0, 6.0, 2.55)
    text_box(s, x=7.00, y=2.10, w=5.7, h=0.4,
             text="COCA-COLA · «HOLIDAYS ARE COMING» AI · 2024",
             size=11, bold=True, color=TEAL)
    add_image(s, str(ASSETS / "screenshots/s30-coca-secret.png"),
              x=7.00, y=2.45, w=5.70, h=2.00)
    # Two разворот тональности chips below
    # Toys R Us swing
    filled_rect(s, 0.55, 4.70, 6.0, 1.60, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.08)
    text_box(s, x=0.70, y=4.78, w=5.7, h=0.30,
             text="TOYS R US · РАЗВОРОТ ТОНАЛЬНОСТИ (июнь 2024)",
             size=10, bold=True, color=TEAL)
    text_box(s, x=0.70, y=5.10, w=2.9, h=0.5,
             text="ПОЛОЖИТЕЛЬНЫЕ", size=10, color=SLATE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=0.70, y=5.45, w=2.9, h=0.5,
             text="+12.2%  →  +3.4%", size=14, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=3.65, y=5.10, w=2.7, h=0.5,
             text="ОТРИЦАТЕЛЬНЫЕ", size=10, color=SLATE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=3.65, y=5.45, w=2.7, h=0.5,
             text="13.5%  →  53.4%", size=14, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=0.70, y=5.90, w=5.7, h=0.30,
             text="−8.8 п.п. положительных · +39.9 п.п. отрицательных", size=10, italic=True, color=DEEP)
    # Coca-Cola swing
    filled_rect(s, 6.85, 4.70, 6.0, 1.60, SURFACE, stroke=LIGHT, stroke_pt=1.0,
                radius=True, radius_adj=0.08)
    text_box(s, x=7.00, y=4.78, w=5.7, h=0.30,
             text="COCA-COLA AI-РЕКЛАМА · НЕГАТИВ (ноя 2024)",
             size=10, bold=True, color=TEAL)
    text_box(s, x=7.00, y=5.10, w=5.7, h=0.5,
             text="«Бездушно» · «жуткие лица» · «колёса грузовика крутятся не туда»",
             size=11, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    text_box(s, x=7.00, y=5.65, w=5.7, h=0.5,
             text="100 сотрудников · 70 000 AI-клипов · вирусные насмешки",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    # Lesson
    lesson_box(s, 0.55, 6.5, 12.3, 0.9,
               "AI-реклама возможна, но эталонная сезонная кампания БЕЗ человеческого лидерства = ущерб бренду. Разворот тональности, а НЕ CTR — главная метрика риска доверия к бренду.")
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """Замещение consolidated 3-stat."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="Замещение: сводная картина — структурный сдвиг (Кейс 11)",
             size=24, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Снижение зарплат снизу — структурное, не временный шок. Оговорки помогают верхнему сегменту, а не нижнему.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 stat blocks horizontal
    stats = [
        ("−17.01%", "работ графического дизайна\nUpwork год к году", MID, "Jobbers Index"),
        ("40% / <10%", "работ $10-19/час vs\n$60+/час · распознан как AI",
         LIGHT, "снижение зарплат снизу"),
        ("4 года", "SAG-AFTRA + WGA\nпродление 2026", GOLD,
         "Оговорка о цифровых репликах"),
    ]
    card_y = 2.2
    card_h = 3.6
    card_w = 3.95
    gap = 0.18
    start_x = 0.55
    for i, (big, label, color, footer) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        # Color band
        filled_rect(s, x, card_y, card_w, 0.6, color, radius=True, radius_adj=0.10)
        # Big stat
        text_box(s, x=x + 0.15, y=card_y + 0.8, w=card_w - 0.30, h=1.4,
                 text=big, size=44, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Label
        text_box(s, x=x + 0.15, y=card_y + 2.3, w=card_w - 0.30, h=0.7,
                 text=label, size=13, color=DEEP, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.30)
        # Footer chip
        chip(s, x + 0.30, card_y + card_h - 0.55, card_w - 0.60, 0.40, footer,
             fill=SURFACE, stroke=LIGHT, color=DEEP, size=10, bold=False)
    # Getty merger chip below
    text_box(s, x=0.55, y=6.05, w=12.3, h=0.4,
             text="+ Слияние Getty + Shutterstock $3.7 млрд (янв 2025) — оборонительная консолидация",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Lesson
    lesson_box(s, 0.55, 6.45, 12.3, 0.75,
               "Замещение — структурное, не временный шок. Оговорки помогают, но снижение зарплат снизу остаётся. Понимай, какой класс работников твой AI-проект вытеснит, ДО запуска.")
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    build_section_divider(p, here_idx=4, title="AI здесь не нужен",
                          frame_phrase="Критерии негативного выбора · 3 зоны только человеком · эмпирическое отторжение конечными пользователями.",
                          notes_slide_id="s32")


def build_s33(p):
    """4 критерия отказа от AI."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="4 критерия отказа от AI в творческом проекте",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Выученные уроки Раздела 3 → переведены в чек-лист.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 2x2 grid of criteria
    grid_x = 0.55
    grid_y = 2.0
    grid_w = 12.3
    grid_h = 4.6
    cell_w = (grid_w - 0.20) / 2
    cell_h = (grid_h - 0.20) / 2
    criteria = [
        ("1. Лицензия на обучающие данные",
         "Нет документированного лицензированного корпуса → юридический долг по категории 1.",
         "→ Andersen против Stability · RIAA против Suno",
         "Firefly = да · Stable Diffusion = риски",
         MID),
        ("2. Проверка сходства результата",
         "Выводы могут воспроизвести охраняемый контент → ответственность по категории 2.",
         "→ теория дословного цитирования NYT",
         "Технический контроль обязателен",
         LIGHT),
        ("3. Согласие на голос/образ",
         "Нет явного согласия → риск класса ScarJo + SAG-AFTRA + Корея.",
         "→ ScarJo против OpenAI · кризис в Корее",
         "Узнаваемые реальные люди — всегда согласие",
         TEAL),
        ("4. Риск доверия к бренду",
         "Эталонная/историческая кампания без человеческого лидерства → измеримая негативная реакция.",
         "→ Coca-Cola · Toys R Us · SI",
         "Разворот тональности, а не CTR",
         GOLD),
    ]
    for i, (title, body, cases, tip, accent) in enumerate(criteria):
        row = i // 2; col = i % 2
        x = grid_x + col * (cell_w + 0.20)
        y = grid_y + row * (cell_h + 0.20)
        ocean_box(s, x, y, cell_w, cell_h)
        filled_rect(s, x, y, cell_w, 0.55, accent, radius=True, radius_adj=0.10)
        accent_text = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=y + 0.08, w=cell_w - 0.30, h=0.40,
                 text=title, size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.18, y=y + 0.70, w=cell_w - 0.36, h=0.60,
                 text=body, size=11, color=DEEP, line_spacing=1.25)
        text_box(s, x=x + 0.18, y=y + 1.40, w=cell_w - 0.36, h=0.35,
                 text=cases, size=10, italic=True, color=GOLD, bold=True,
                 line_spacing=1.15)
        # Tip chip — separate solid block to avoid overlap
        chip(s, x + 0.18, y + cell_h - 0.50, cell_w - 0.36, 0.38, tip,
             fill=TEAL_TINT, stroke=TEAL, color=DEEP, size=10, bold=True)
    speaker_notes(s, load_notes("s33"))


def build_s34(p):
    """3 zones только человеком."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="3 зоны, где AI не должен подменять человека",
             size=28, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="«AI как инструмент» работает; «AI как замещение» — нет.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 columns
    card_y = 2.0
    card_h = 4.6
    card_w = (12.3 - 0.30) / 3
    start_x = 0.55
    gap = 0.15
    zones = [
        ("Расследовательская журналистика",
         "NYT/WaPo гайдлайны запрещают AI для первичного репортажа",
         [
             "Проверка источников — только человеком",
             "Интервью под запись — только человеком",
             "Журналистика подотчётности — человеческая ответственность не делегируема",
         ], MID),
        ("Оригинальное творческое руководство",
         "Эталонные брендовые кампании без человеческого лидерства = ущерб бренду",
         [
             "Toys R Us Cannes 2024 — разворот тональности −10 п.п.",
             "Coca-Cola Christmas 2024 — прохладный приём",
             "SI by-line — расследовательский скандал",
         ], LIGHT),
        ("Связный нарратив длинного формата",
         "Альбом 50 мин · многоактный сценарий — пока человек",
         [
             "Дрейф непрерывности персонажа после 5-10 сцен",
             "Связность сюжета через 90+ минут — нет",
             "Связность голоса через целый альбом — нет",
         ], GOLD),
    ]
    for i, (ztitle, zsubtitle, zpoints, accent) in enumerate(zones):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        filled_rect(s, x, card_y, card_w, 0.6, accent, radius=True, radius_adj=0.10)
        accent_text = WHITE if accent != GOLD else DEEP
        text_box(s, x=x + 0.15, y=card_y + 0.08, w=card_w - 0.30, h=0.45,
                 text=ztitle, size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=x + 0.18, y=card_y + 0.80, w=card_w - 0.36, h=0.8,
                 text=zsubtitle, size=12, italic=True, color=MID, line_spacing=1.30)
        for j, point in enumerate(zpoints):
            y = card_y + 1.80 + j * 0.85
            filled_rect(s, x + 0.18, y, 0.10, 0.10, accent, radius=True,
                        radius_adj=0.5)
            text_box(s, x=x + 0.40, y=y - 0.05, w=card_w - 0.55, h=0.75,
                     text=point, size=11, color=DEEP, line_spacing=1.35)
    # Footer
    text_box(s, x=0.55, y=6.85, w=12.3, h=0.4,
             text="Здесь — AI только как вспомогательный инструмент под надзором человека; никогда как слой замещения.",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s34"))


def build_s35(p):
    """YouTube AI-превью."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="AI-превью на YouTube · 47.3% креаторов отказались",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Эмпирическое отторжение конечными пользователями · опрос Social Blade Creator Survey · Дек 2025",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 3 stat blocks
    stats = [
        ("47.3%", "креаторов отказались от\nAI-превью", GOLD),
        ("−22% / −19%", "падение CTR:\nжуткая кожа / провал текста на мобильных", LIGHT),
        ("−61.8%", "отвал в первые 15 сек\nнесоответствие обещания и контента", MID),
    ]
    card_y = 2.0
    card_h = 3.0
    card_w = 3.95
    gap = 0.18
    start_x = 0.55
    for i, (big, label, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        ocean_box(s, x, card_y, card_w, card_h)
        filled_rect(s, x, card_y, card_w, 0.6, color, radius=True, radius_adj=0.10)
        accent_text = WHITE if color != GOLD else DEEP
        text_box(s, x=x + 0.15, y=card_y + 0.08, w=card_w - 0.30, h=0.45,
                 text=f"#{i+1}", size=15, bold=True, color=accent_text,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=x + 0.15, y=card_y + 0.85, w=card_w - 0.30, h=1.0,
                 text=big, size=36, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=x + 0.20, y=card_y + 1.95, w=card_w - 0.40, h=1.0,
                 text=label, size=12, color=DEEP, italic=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.40)
    # Lesson (expanded to full width, fills freed space below 3 stat cards)
    lesson_box(s, 0.55, 5.30, 12.3, 1.85,
               "Конечные пользователи замечают и наказывают AI-срезание углов. Риск доверия к бренду не теоретический — измеримые CTR + отвал + удержание. Если продукт опирается на AI-визуалы для конечных пользователей, измеряй не только стоимость генерации — измеряй CTR, удержание, отношение к бренду.")
    speaker_notes(s, load_notes("s35"))


def build_s36(p):
    build_section_divider(p, here_idx=5, title="Что инженеру делать",
                          frame_phrase="Чек-лист к действию · 5-вопросный · главный итог лекции.",
                          notes_slide_id="s36")


def build_s37(p):
    """5-question checklist."""
    s = blank(p)
    text_box(s, x=0.55, y=0.45, w=12.3, h=0.75,
             text="5-вопросный чек-лист перед AI в творческом проекте",
             size=26, bold=True, color=DEEP, align=PP_ALIGN.LEFT)
    text_box(s, x=0.55, y=1.20, w=12.3, h=0.5,
             text="Применяй ДО старта, не после. Если хоть один «нет/риск» — пересмотри подход.",
             size=14, italic=True, color=MID, align=PP_ALIGN.LEFT)
    # 5 questions vertically
    questions = [
        ("1.", "Лицензирование обучающих данных у инструмента?",
         "Firefly = да (Adobe Stock + лицензированный контент) · SD/Midjourney = риски класса Andersen"),
        ("2.", "Проверка сходства результата с охраняемым контентом?",
         "Риск дословного цитирования NYT · технический контроль обязателен"),
        ("3.", "Согласие на голос/образ — если применимо?",
         "ScarJo · SAG-AFTRA Цифровые реплики · Корея — всегда согласие для реальных людей"),
        ("4.", "IP-чистые инструменты для коммерческого применения?",
         "Сквозной pipeline: лицензированный корпус + проверка сходства + раскрытие"),
        ("5.", "Риск доверия к бренду для эталонных/исторических кампаний?",
         "Coca-Cola · Toys R Us · SI — измеримая трата капитала бренда"),
    ]
    q_x = 0.55
    q_w = 8.5
    q_y = 2.0
    q_h = 0.95
    q_gap = 0.10
    for i, (num, qtitle, qbody) in enumerate(questions):
        y = q_y + i * (q_h + q_gap)
        ocean_box(s, q_x, y, q_w, q_h)
        # Number badge
        filled_rect(s, q_x + 0.10, y + 0.12, 0.7, q_h - 0.24, MID,
                    radius=True, radius_adj=0.20)
        text_box(s, x=q_x + 0.10, y=y + 0.12, w=0.7, h=q_h - 0.24,
                 text=num, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=q_x + 0.95, y=y + 0.10, w=q_w - 1.05, h=0.40,
                 text=qtitle, size=14, bold=True, color=DEEP, line_spacing=1.0)
        text_box(s, x=q_x + 0.95, y=y + 0.50, w=q_w - 1.05, h=0.40,
                 text=qbody, size=11, italic=True, color=SLATE, line_spacing=1.20)
    # Right: decision branch
    ocean_box(s, 9.25, 2.0, 3.7, 5.25, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.8)
    text_box(s, x=9.45, y=2.15, w=3.4, h=0.4,
             text="ЕСЛИ «НЕТ/РИСК» — 3 ВАРИАНТА",
             size=11, bold=True, color=GOLD)
    options = [
        ("A.", "Не использовать AI", "Альтернатива без AI", MID),
        ("B.", "Структурно митигировать",
         "Лицензирование + проверка сходства + инфраструктура согласия + учёт бренда", LIGHT),
        ("C.", "Явно принять риск",
         "Документированное бизнес-решение + выверенная митигация. НЕ неявно.", TEAL),
    ]
    for i, (oletter, ot, ob, oc) in enumerate(options):
        y = 2.65 + i * 1.50
        filled_rect(s, 9.45, y, 0.5, 0.5, oc, radius=True, radius_adj=0.20)
        text_box(s, x=9.45, y=y, w=0.5, h=0.5, text=oletter,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=10.05, y=y, w=2.85, h=0.45,
                 text=ot, size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=9.45, y=y + 0.55, w=3.4, h=0.85,
                 text=ob, size=11, italic=True, color=SLATE, line_spacing=1.30)
    speaker_notes(s, load_notes("s37"))


def build_s38(p):
    """Q&A."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # Left: big Q&A
    text_box(s, x=0.55, y=1.5, w=5.5, h=4.5, text="Q&A?",
             size=200, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # Right: central question recap + резервные вопросы
    ocean_box(s, 6.30, 1.5, 6.65, 5.0)
    text_box(s, x=6.50, y=1.65, w=6.3, h=0.4,
             text="ЦЕНТРАЛЬНЫЙ ВОПРОС — НАПОМИНАНИЕ",
             size=11, bold=True, color=TEAL)
    text_box(s, x=6.50, y=2.10, w=6.3, h=1.4,
             text="Что AI сделал с креативной индустрией к 2026 — и где сказать «нет»?",
             size=20, italic=True, bold=True, color=DEEP, line_spacing=1.30)
    text_box(s, x=6.50, y=3.65, w=6.3, h=0.4,
             text="РЕЗЕРВНЫЕ ВОПРОСЫ",
             size=11, bold=True, color=TEAL)
    prompts = [
        "Где границы «добросовестного использования» при обучении AI?",
        "Sora vs Lionsgate — где Голливуд в 2030?",
        "Минцифры законопроект — что меняется для RU-инженеров?",
    ]
    for i, prompt in enumerate(prompts):
        y = 4.10 + i * 0.65
        filled_rect(s, 6.50, y, 6.3, 0.55, SURFACE,
                    stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.20)
        text_box(s, x=6.65, y=y, w=6.1, h=0.55,
                 text=f"·  {prompt}", size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s38"))


def build_s39(p):
    """Closing с hero-image (X-62 VISTA) — bridge к Лекции 9 (AI в авиакосмосе и обороне)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    # ---- LEFT: thank you (compacted) ----
    ocean_box(s, 0.45, 1.30, 5.40, 4.95)
    text_box(s, x=0.65, y=1.50, w=5.05, h=0.45,
             text="СПАСИБО", size=18, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT)
    text_box(s, x=0.65, y=2.15, w=5.05, h=1.50,
             text="за внимание",
             size=40, bold=True, color=DEEP, line_spacing=1.10)
    text_box(s, x=0.65, y=3.95, w=5.05, h=0.45,
             text="Чек-лист и источники — по QR ниже",
             size=13, italic=True, color=MID, line_spacing=1.20)
    # QR placeholder
    filled_rect(s, 0.65, 4.75, 1.0, 1.0, WHITE, stroke=DEEP, stroke_pt=1.5)
    text_box(s, x=0.65, y=4.75, w=1.0, h=1.0, text="QR",
             size=18, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, x=1.80, y=4.75, w=3.90, h=1.0,
             text="Полная библиография источников лекции (chapter.md + research-досье)",
             size=10, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.30)
    # ---- RIGHT: HERO BRIDGE IMAGE (X-62 VISTA — Лекция 9 foreshadow) ----
    # Image area: 7.0×3.95 (preserving 16:9 aspect of X-62 photo) ≈ 27.7 sq in ≈ 28%
    # Combined with right text panel — image+caption block ≈ 40% площади
    ocean_box(s, 6.10, 0.55, 6.75, 4.30, fill=WHITE, stroke=LIGHT, stroke_pt=1.5)
    add_image(s, str(ASSETS / "screenshots/s39-x62-vista.jpg"),
              x=6.20, y=0.65, w=6.55, h=4.10)
    # Attribution chip under image
    text_box(s, x=6.20, y=4.90, w=6.55, h=0.30,
             text="X-62 VISTA · USAF Test Pilot School · DARPA ACE AI dogfight 2023 · Wikimedia Commons",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)
    # Next lecture text block under hero
    text_box(s, x=6.20, y=5.30, w=6.55, h=0.35,
             text="СЛЕДУЮЩАЯ ЛЕКЦИЯ", size=12, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT)
    text_box(s, x=6.20, y=5.65, w=6.55, h=0.75,
             text="AI в авиакосмической отрасли и оборонном комплексе",
             size=20, bold=True, color=DEEP, line_spacing=1.15, align=PP_ALIGN.LEFT)
    filled_rect(s, 6.20, 6.50, 0.05, 0.45, fill=GOLD)
    text_box(s, x=6.35, y=6.45, w=6.50, h=0.55,
             text="От публичной контактной поверхности — к безопасности-критичному",
             size=13, italic=True, color=MID, align=PP_ALIGN.LEFT, line_spacing=1.20)
    speaker_notes(s, load_notes("s39"))


# ============================================================
# Main
# ============================================================

def main():
    p = setup_pres()
    # Сборка в порядке deck.yaml
    builders = [
        build_s01, build_s02, build_s03, build_s04, build_s05, build_s05a,
        build_s06, build_s07, build_s08, build_s09, build_s10, build_s10a,
        build_s11,
        build_s13, build_s14, build_s15, build_s16, build_s17,
        build_s19, build_s20, build_s21, build_s22, build_s23, build_s24,
        build_s25, build_s26, build_s27, build_s28, build_s29, build_s30,
        build_s31,
        build_s32, build_s33, build_s34, build_s35,
        build_s36, build_s37,
        build_s38, build_s39,
    ]
    for builder in builders:
        builder(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved {OUT} with {len(builders)} slides.")


if __name__ == "__main__":
    main()
