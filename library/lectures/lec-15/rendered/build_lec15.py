"""
Full 39-slide build of Лекции 15 «AI в научных исследованиях».

Source-of-truth: deck.yaml v1 + chapter v2.2 multi-part (~32 850 слов) + slides/*.md.

Issue #143 · Phase 6 visual rendering · downstream от chapter (book-first).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml ≈ 75 мин.

Lec-N-1 pattern compliance: match lec-14 (cover + lecture-map + 5 section dividers +
dedicated Q&A; top progress bar только на dividers + cover).

Build via: python3 build_lec15.py — generates lec-15.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
RED_WARN = RGBColor(0xC0, 0x39, 0x2B)
ROADMAP = RGBColor(0xD9, 0xE2, 0xEC)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "lec-15.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *, size=16, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except: pass
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline_box(slide, x, y, w, h, lines, *, size=14, bold=False, color=DEEP,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
                  line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except: pass
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, *, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def down_arrow(slide, x, y, w, h, *, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w, h):
    if not path or not Path(path).exists():
        return None
    try:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                         width=Inches(w), height=Inches(h))
    except Exception as e:
        print(f"Image fail {path}: {e}")
        return None


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def roadmap_bar(slide, current_section):
    """6-section roadmap bar at top of section dividers + cover."""
    sections = ["Введение", "Hypothesis+Design", "Experiment", "Analyse", "Write+Review", "Замыкание"]
    bar_y = 0.4
    bar_h = 0.32
    total_w = 12.33
    seg_w = total_w / 6
    for i, name in enumerate(sections):
        x = 0.5 + i * seg_w
        is_active = (i == current_section)
        rectangle(slide, x, bar_y, seg_w - 0.05, bar_h,
                  fill=GOLD if is_active else ROADMAP)
        text_box(slide, x, bar_y + 0.04, seg_w - 0.05, bar_h - 0.04,
                 f"{i+1}. {name}", size=10, bold=is_active,
                 color=DEEP if is_active else SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, text, *, y=7.05):
    text_box(slide, 0.5, y, 12.33, 0.35, text,
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)


def attribution(slide, text, *, x=0.5, y=6.95, w=12.33):
    text_box(slide, x, y, w, 0.3, text,
             size=10, italic=True, color=SLATE, align=PP_ALIGN.LEFT)


def slide_header(slide, assertion, *, y=0.4):
    """Standard assertion header for content slides."""
    text_box(slide, 0.5, y, 12.33, 0.9, assertion,
             size=24, bold=True, color=DEEP, line_spacing=1.1)


# ========== INDIVIDUAL SLIDE BUILDERS ==========
# Imported from build_lec15_slides.py
