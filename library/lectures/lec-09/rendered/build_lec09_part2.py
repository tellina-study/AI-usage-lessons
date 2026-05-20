"""
Part 2 of build_lec09.py — slides 9 through 43.
Imports helpers and constants from build_lec09 module.
"""
from pathlib import Path
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Re-import constants and helpers from main builder
from build_lec09 import (
    DEEP, MID, LIGHT, TEAL, SURFACE, WHITE, GOLD, SLATE,
    COVER_OUTLINE, GOLD_TINT, TEAL_TINT, LIGHT_TINT, SOFT_GREY, DARK_GREY, RED_WARN,
    FONT_HEAD, FONT_BODY, SLIDE_W_IN, SLIDE_H_IN, ASSETS,
    blank, set_slide_bg, disable_shadow,
    text_box, text_runs, ocean_box, filled_rect, hr_line, add_arrow,
    add_image, add_speaker_notes,
    add_progress_bar, add_footer, add_assertion_title,
    section_divider,
)


# ========== Generic builders ==========

def assertion_with_cards(prs, assertion, cards_data, *, footer=None, notes="",
                          title_size=24, card_layout="row"):
    """Generic card-layout slide. cards_data: [(icon, title, body_lines, accent_color)]"""
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, assertion, size=title_size)

    n = len(cards_data)
    if card_layout == "row":
        # All cards in single row
        gap = 0.25
        total_w = 12.13
        card_w = (total_w - gap * (n - 1)) / n
        card_h = 4.5
        for i, item in enumerate(cards_data):
            icon, title, body_lines, accent = item
            x = 0.6 + i * (card_w + gap)
            y = 2.0
            ocean_box(s, x, y, card_w, card_h)
            # Icon
            icon_p = ASSETS / "icons" / f"{icon}-96.png"
            if icon_p.exists():
                add_image(s, icon_p, x + (card_w - 0.8) / 2, y + 0.3, w=0.8, h=0.8)
            text_box(s, x + 0.2, y + 1.25, card_w - 0.4, 0.6, title,
                     size=18, bold=True, color=accent,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            ly = y + 1.95
            for line in body_lines:
                text_box(s, x + 0.25, ly, card_w - 0.5, 0.5,
                         "• " + line, size=11, color=MID, line_spacing=1.25)
                ly += 0.45
    elif card_layout == "2x2":
        gap = 0.25
        card_w = (12.13 - gap) / 2; card_h = 2.2
        for i, item in enumerate(cards_data):
            icon, title, body_lines, accent = item
            row = i // 2; col = i % 2
            x = 0.6 + col * (card_w + gap)
            y = 2.0 + row * (card_h + gap)
            ocean_box(s, x, y, card_w, card_h)
            icon_p = ASSETS / "icons" / f"{icon}-48.png"
            if icon_p.exists():
                add_image(s, icon_p, x + 0.2, y + 0.25, w=0.55, h=0.55)
            text_box(s, x + 0.95, y + 0.25, card_w - 1.1, 0.5, title,
                     size=16, bold=True, color=accent)
            ly = y + 0.85
            for line in body_lines:
                text_box(s, x + 0.3, ly, card_w - 0.5, 0.35,
                         "• " + line, size=11, color=MID, line_spacing=1.2)
                ly += 0.35

    if footer:
        add_footer(s, footer)
    if notes:
        add_speaker_notes(s, notes)
    return s


# ========== Concrete slide builders ==========

def slide_09_constellation(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "BlackSky + Planet + Capella + ICEYE — четыре комплементарных подхода")

    # Table-style 4 rows
    headers = ["Игрок", "Подход", "Контракт", "Тип"]
    rows = [
        ("BlackSky Gen-3 + Spectra AI", "Малые сат, частая revisit",
         "$100M+ subscription", "EO + CV"),
        ("Planet Labs Dove", "Сотни сат, ежедневное покрытие",
         "NRO EOCL $146M+", "EO"),
        ("Capella Space", "SAR — всепогодный",
         "NGA partnership", "SAR"),
        ("ICEYE (Финляндия)", "SAR коммерческий",
         "Украина + НАТО", "SAR"),
    ]
    table_x = 0.6; table_w = 8.5
    col_widths = [3.0, 2.8, 1.6, 1.1]
    cum_x = 0
    cell_h = 0.8
    header_y = 2.0

    # Header row
    for i, h in enumerate(headers):
        cx = table_x + cum_x
        filled_rect(s, cx, header_y, col_widths[i], 0.5, MID)
        text_box(s, cx + 0.1, header_y + 0.05, col_widths[i] - 0.2, 0.4, h,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cum_x += col_widths[i]

    # Body rows
    body_y = header_y + 0.5
    for r_i, row in enumerate(rows):
        bg = SURFACE if r_i % 2 == 0 else WHITE
        cum_x = 0
        for c_i, cell in enumerate(row):
            cx = table_x + cum_x
            filled_rect(s, cx, body_y + r_i * cell_h, col_widths[c_i], cell_h, bg,
                        stroke=LIGHT, stroke_pt=0.5)
            text_box(s, cx + 0.12, body_y + r_i * cell_h + 0.1,
                     col_widths[c_i] - 0.24, cell_h - 0.2, cell,
                     size=11, color=DEEP if c_i == 0 else MID,
                     bold=(c_i == 0), line_spacing=1.2,
                     anchor=MSO_ANCHOR.MIDDLE)
            cum_x += col_widths[c_i]

    # Right — constellation diagram
    ocean_box(s, 9.4, 2.0, 3.4, 3.7)
    text_box(s, 9.55, 2.15, 3.1, 0.35, "Орбитальный микс",
             size=13, bold=True, color=MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # Earth circle in middle
    earth = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.55), Inches(3.55),
                               Inches(1.1), Inches(1.1))
    earth.fill.solid(); earth.fill.fore_color.rgb = SURFACE
    earth.line.color.rgb = LIGHT; earth.line.width = Pt(1.5)
    text_box(s, 10.55, 3.55, 1.1, 1.1, "Earth",
             size=10, italic=True, color=MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 4 colored dots around
    dots = [(0.5, 2.7, MID), (1.6, 2.55, LIGHT), (0.4, 4.7, TEAL), (1.7, 4.85, GOLD)]
    for dx, dy, col in dots:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(9.4 + dx), Inches(2.0 + dy),
                                 Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background()

    # Bottom callout
    ocean_box(s, 0.6, 6.4, 12.13, 0.6, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.48, 11.63, 0.45,
             "Главное преимущество SAR — видит сквозь облака и ночью. ICEYE активен в Украине с 2022.",
             size=13, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    add_speaker_notes(s,
        "Четвёрка коммерческих игроков формирует рынок 2026 года, и у каждого "
        "свой стратегический подход. BlackSky сделал ставку на малые спутники "
        "с высокой частотой revisit. Их Gen-3 спутники работают со Spectra AI, "
        "который делает CNN-классификацию и change detection. Planet Labs — "
        "самая массовая сеть с сотнями малых спутников Dove и ежедневным "
        "глобальным покрытием. Главный контракт — NRO EOCL. Capella Space и "
        "ICEYE — SAR-операторы. Преимущество SAR — всепогодная съёмка, "
        "способность видеть сквозь облака и ночью. Главное применение — "
        "морское наблюдение, отслеживание «теневого флота», мониторинг военной "
        "инфраструктуры. Capella сотрудничает с NGA, ICEYE — финский оператор с "
        "активными контрактами в Украине и НАТО. Структурный takeaway: вы не "
        "выбираете «лучшего» вендора. Вы выбираете правильный сенсор и "
        "правильную частоту покрытия для вашей задачи. AI работает поверх этого "
        "выбора, а не вместо него.")
    return s


def slide_10_edge_ai(prs):
    cards = [
        ("flask-conical", "Demonstrators",
         ["ESA Φ-sat-2 (август 2024)", "Planetek AI-eXpress на Jetson Orin NX"],
         MID),
        ("activity", "Production telemetry",
         ["Lockheed Pony Express 2 + T-TAURI", "Onboard ML аномалии телеметрии"],
         MID),
        ("radar", "SDA tracking",
         ["Slingshot Agatha + TALOS · июль 2025", "204 сенсора · 21 локация · 5 континентов"],
         TEAL),
        ("database", "Commercial archive",
         ["TerraTech / Роскосмос", "Гражданская onboard-классификация"],
         LIGHT),
    ]
    s = assertion_with_cards(prs,
        "Edge AI on-orbit: ML на спутнике, не на земле. ESA Φ-sat-2 — remote-upgradable.",
        cards, card_layout="2x2",
        footer="Adoption: от «AI-derived detection in hours» к «predictive intelligence before event». Сценарий 2026-2028.",
        notes=(
            "Параллельно с большой коммерческой аналитикой развивается отдельная "
            "линия — edge AI on-orbit. Это ML-вычисления прямо на спутнике, без "
            "передачи сырого сигнала на землю. Цель — снизить латентность и "
            "ширину канала: вместо мегабайтов сырого изображения передаётся "
            "килобайт сводки «вот тут изменение».\n\n"
            "Программы 2024-2026 года группируются в четыре категории. "
            "Демонстраторы: ESA Φ-sat-2 — европейский demonstration satellite с "
            "remote-upgradable ML-моделями. После запуска модель можно дообучить "
            "и заменить новым весом по телекомандам, без замены оборудования. "
            "AI-eXpress 1+ от Planetek Italia — серия европейских edge-computing "
            "спутников на NVIDIA Jetson Orin NX.\n\n"
            "Production telemetry: Lockheed Pony Express 2 + T-TAURI — "
            "американский военный production-аналог. Это уже не demo, это "
            "инфраструктура. SDA tracking: Slingshot Agatha + TALOS — 204 "
            "сенсора в 21 локации на 5 континентах плюс ML-«отпечатки» спутников "
            "по фотометрическому паттерну. Это не ISR в смысле наблюдения за "
            "землёй, это наблюдение за самим космосом.\n\n"
            "Commercial archive: TerraTech от Роскосмоса — гражданская "
            "коммерческая edge-аналитика. Adoption-направление по Sense — "
            "растёт быстро."))
    return s


def slide_11_russian_sat(prs):
    cards = [
        ("globe", "ТЕРРА ТЕХ",
         ["Дочерняя структура Роскосмоса (2017)",
          "BRICS agriculture monitoring (2024)",
          "Spatial data + ML-классификация",
          "Caveat: объём метрик не публикуется"],
         GOLD),
        ("satellite-dish", "СКАНЭКС",
         ["Единственная direct-receiving в РФ/СНГ",
          "3,5М+ архив снимков",
          "Эксклюзив для Яндекс.Карт",
          "Court ban на >2м distribution"],
         MID),
        ("satellite", "СПУТНИКС",
         ["100+ кубсатов с 2013",
          "Zorkiy-2M · 2,5м · 4 спектр канала",
          "Sitronics 45 cubsats 2024 (53% RU)"],
         LIGHT),
    ]
    s = assertion_with_cards(prs,
        "Российский слой Sense — ТЕРРА ТЕХ, СКАНЭКС, СПУТНИКС",
        cards, card_layout="row",
        footer="Оборонный Russian Sense — открытых данных мало. Geran-2 — Раздел 3; российский C2 — Раздел 2.",
        notes=(
            "В российском контексте Sense развивается тремя путями, и публичная "
            "часть из этих путей — гражданская спутниковая аналитика.\n\n"
            "ТЕРРА ТЕХ — дочерняя структура Роскосмоса, основана в 2017 году. "
            "Предоставляет digital services на пространственных данных. Главный "
            "публичный кейс — мониторинг сельскохозяйственных земель в странах "
            "БРИКС по соглашению 2024 года. Точный объём и эффективность не "
            "публикуются — это не PR в смысле «ничего нет», но и не полная "
            "transparency.\n\n"
            "СКАНЭКС — единственная компания в РФ и СНГ, напрямую принимающая "
            "данные дистанционного зондирования Земли на собственные наземные "
            "станции. Архив — три с половиной миллиона снимков; эксклюзивный "
            "поставщик для Яндекс.Карт. Московский суд ограничил распространение "
            "снимков разрешением выше 2 метров — security restriction 2014 года.\n\n"
            "СПУТНИКС — более 100 кубсатов с 2013 года. Главная констелляция — "
            "Зоркий-2М: три спутника на орбите плюс ещё три, мультиспектральная "
            "камера 2,5 метра и 4 спектральных канала. Sitronics Group развернул "
            "в 2024 году 45 коммерческих кубсатов — 53% всех российских "
            "развёртываний этого года.\n\n"
            "В оборонном Sense российских открытых данных мало. Российский "
            "гражданский слой Sense ассимметричен западному: вендоров меньше, "
            "объём публикуемых метрик меньше, но реальные операционные кейсы "
            "есть."))
    return s


def slide_12_predictive_maintenance(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Predictive maintenance — массовая гражданская AI-инфраструктура")

    # Left: Skywise chart
    ocean_box(s, 0.6, 2.0, 6.5, 4.4)
    chart_p = ASSETS / "charts" / "c12-skywise-bar.png"
    if chart_p.exists():
        add_image(s, chart_p, 0.75, 2.2, w=6.2, h=2.9)
    text_box(s, 0.85, 5.3, 6.0, 0.95,
             "easyJet: 8,1 тонны топлива/ВС/год сэкономлено; "
             "44 предотвращённые отмены рейсов (июль 2024)",
             size=12, italic=True, color=MID, line_spacing=1.25)

    # Right: Rolls-Royce + stack cards
    ocean_box(s, 7.4, 2.0, 5.4, 2.2)
    text_box(s, 7.6, 2.15, 5.0, 0.45, "Rolls-Royce IntelligentEngine",
             size=16, bold=True, color=DEEP)
    text_box(s, 7.6, 2.65, 5.0, 0.35,
             "С 2018 · digital twin каждого двигателя",
             size=11, italic=True, color=MID)
    text_runs(s, 7.6, 3.1, 5.0, 1.05, [
        {"text": "~400", "size": 36, "bold": True, "color": GOLD},
        {"text": " предотвращённых событий", "size": 14, "color": DEEP},
        {"newpara": True, "text": "обслуживания в год", "size": 12, "italic": True, "color": MID},
    ])

    ocean_box(s, 7.4, 4.4, 5.4, 1.95, fill=LIGHT_TINT)
    text_box(s, 7.6, 4.55, 5.0, 0.4, "Стек публичный",
             size=14, bold=True, color=DEEP)
    stack = ["Microsoft Azure data lake", "Databricks lakehouse", "ML pipelines"]
    for i, item in enumerate(stack):
        text_box(s, 7.7, 4.95 + i * 0.4, 5.0, 0.35, "• " + item,
                 size=12, color=MID)

    add_footer(s, "В обороне аналог — F-35 ALIS → ODIN. См. следующий слайд")
    add_speaker_notes(s,
        "Помимо разведки, AI в Sense массово работает на собственных аппаратах. "
        "Predictive maintenance — это семейство задач: по телеметрии двигателей, "
        "систем, бортового оборудования предсказать отказ компонента до того, "
        "как он случится, и заменить его на плановом обслуживании, а не на "
        "аварийной посадке.\n\n"
        "Rolls-Royce IntelligentEngine плюс TotalCare работает с 2018 года и "
        "сейчас представляет собой digital twin каждого летающего двигателя "
        "плюс ML-конвейеры на телеметрии. Стек публичный: Microsoft Azure data "
        "lake, поверх него Databricks lakehouse, далее ML pipelines. Главная "
        "метрика — около 400 непланированных событий обслуживания "
        "предотвращаются в год.\n\n"
        "Airbus Skywise — более широкая платформа. К концу 2024 года к "
        "платформе подключены около 11 600 самолётов; около 40 авиакомпаний "
        "на расширенной подписке SFP+, что покрывает около 1 500 ВС. easyJet "
        "с использованием Skywise сообщил об экономии топлива около 8,1 тонны "
        "на ВС в год и о 44 предотвращённых отменах рейсов.\n\n"
        "Что важно для нашего инженерного слуха: это не пилот, это не "
        "лаборатория, это рабочая инфраструктура, которая обслуживает каждый "
        "день тысячи коммерческих самолётов. В оборонном секторе аналог — "
        "F-35 ALIS, и его преемник ODIN. О том, почему ALIS не сработал — "
        "следующий слайд.")
    return s


def slide_13_f35_alis(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "F-35 ALIS нарушил все три условия predictive maintenance — списан в июне 2024",
        size=22)

    # Left: 3 condition cards in vertical stack
    conditions = [
        ("Быстрый feedback loop", "Drift детектируется в годы, не дни. Модель устаревает раньше, чем ошибки видны."),
        ("Доступная ground truth", "Нет способа верифицировать каждый alert. Ложные тревоги накапливаются → доверие падает."),
        ("FP-cost ≤ FN-cost", "Adversarial UX — персонал обходит через Excel. High false-positive — экипаж теряет время на инспекции."),
    ]
    cc_w = 6.8; cc_h = 1.35; gap = 0.15
    for i, (title, desc) in enumerate(conditions):
        y = 1.95 + i * (cc_h + gap)
        ocean_box(s, 0.6, y, cc_w, cc_h, fill=GOLD_TINT, stroke=GOLD)
        # X marker red
        text_box(s, 0.75, y + 0.4, 0.55, 0.6, "✗",
                 size=32, bold=True, color=RED_WARN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, 1.35, y + 0.18, cc_w - 1.4, 0.45, f"{i+1}. {title}",
                 size=16, bold=True, color=DEEP)
        text_box(s, 1.35, y + 0.65, cc_w - 1.4, 0.65, desc,
                 size=12, color=MID, italic=True, line_spacing=1.25)

    # Right: cost chart
    ocean_box(s, 7.6, 1.95, 5.15, 3.0)
    chart_p = ASSETS / "charts" / "c13-f35-cost.png"
    if chart_p.exists():
        add_image(s, chart_p, 7.7, 2.0, w=4.95, h=2.5)
    text_box(s, 7.7, 4.45, 4.95, 0.5,
             "F-35 дороже более сложного F-22 — индикатор системной проблемы",
             size=10, italic=True, color=MID, line_spacing=1.2)

    # ODIN callout
    ocean_box(s, 7.6, 5.05, 5.15, 1.5, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 7.8, 5.15, 4.85, 0.4, "ODIN — другой подход",
             size=14, bold=True, color=DEEP)
    odin_pts = ["Government-owned (no vendor lock)",
                "Меньший охват",
                "Явный HITL для flight-authorisation",
                "Disconnected mode"]
    for i, p in enumerate(odin_pts):
        text_box(s, 7.85, 5.55 + i * 0.22, 4.85, 0.22, "• " + p,
                 size=10, color=MID, line_spacing=1.1)

    add_footer(s, "Источники: GAO-20-316 · GAO-22-105128 · Air & Space Forces 2024")
    add_speaker_notes(s,
        "F-35 ALIS — система предиктивного обслуживания истребителя F-35. К концу "
        "2010-х превратился в источник проблем. Высокая ложноположительная "
        "активность: ALIS помечал самолёт как «no-fly» когда никакой реальной "
        "проблемы не было. Неточные данные: GAO в отчёте 2020 года сообщал — "
        "«Inaccurate and missing data have at times resulted in the system "
        "signalling that an F-35 should not be flown — even though aircraft "
        "had no issues». Adversarial UX: пользоваться ALIS было настолько сложно, "
        "что персонал систематически обходил систему. Cost-per-flight-hour к "
        "пику — $42-44 тысячи, выше чем у F-22.\n\n"
        "Финальная версия ALIS — июнь 2024, начался переход на ODIN — "
        "government-owned, с явно отделённой flight-clearance authority, с "
        "disconnected mode.\n\n"
        "Урок номер один. Predictive maintenance в безопасностно-критичной "
        "области работает только при выполнении трёх условий: быстрый feedback "
        "loop, доступная ground truth, FP-cost меньше или равно FN-cost. ALIS "
        "нарушил все три. ODIN строится в явной осведомлённости об этом "
        "нарушении. Урок номер два. Тот же триплет применим к любой predictive-"
        "системе — от мониторинга турбин до диагностики промышленных линий.")
    return s


def slide_14_adversarial_gps(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Adversarial SAR ATR + GPS-spoofing — accuracy обманчив, single-source хрупок",
        size=22)

    # Left: Adversarial SAR
    ocean_box(s, 0.6, 1.95, 6.0, 4.7)
    text_box(s, 0.8, 2.1, 5.6, 0.4, "Adversarial SAR ATR",
             size=18, bold=True, color=DEEP)
    icon_p = ASSETS / "icons" / "triangle-alert-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 5.9, 2.1, w=0.55, h=0.55)

    text_box(s, 0.8, 2.6, 5.6, 0.5, "Дешёвые металлические рассеиватели в специальной геометрии обманывают classifier",
             size=12, italic=True, color=MID, line_spacing=1.25)

    # Schematic illustration: tank → SAR → wrong label
    for i, (label, x_off) in enumerate([("Танк +\nрассеиватели", 0.0),
                                          ("SAR\nclassifier", 1.95),
                                          ("Wrong\nlabel", 3.9)]):
        bx = 0.85 + x_off
        ocean_box(s, bx, 3.3, 1.7, 1.0, fill=LIGHT_TINT)
        text_box(s, bx, 3.45, 1.7, 0.7, label,
                 size=11, color=DEEP, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        if i < 2:
            add_arrow(s, bx + 1.75, 3.65, 0.15, 0.3, fill=LIGHT)

    text_box(s, 0.8, 4.5, 5.6, 0.35, "Защита:",
             size=13, bold=True, color=DEEP)
    defenses = [
        "Bayesian uncertainty estimates",
        "Adversarial training",
        "Abstention pathway → human"
    ]
    for i, d in enumerate(defenses):
        text_box(s, 0.95, 4.85 + i * 0.3, 5.4, 0.28, "• " + d,
                 size=11, color=MID)

    text_box(s, 0.8, 6.2, 5.6, 0.3, "Source: Du et al. 2024 (arXiv:2312.02912)",
             size=10, italic=True, color=LIGHT)

    # Right: GPS spoofing
    ocean_box(s, 6.8, 1.95, 5.95, 4.7)
    text_box(s, 7.0, 2.1, 5.5, 0.4, "GPS-spoofing civil aviation",
             size=18, bold=True, color=DEEP)
    icon_p = ASSETS / "icons" / "radio-tower-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 12.0, 2.1, w=0.55, h=0.55)

    chart_p = ASSETS / "charts" / "c14-gps-spoof.png"
    if chart_p.exists():
        add_image(s, chart_p, 7.0, 2.65, w=3.5, h=2.0)
    text_box(s, 10.6, 2.95, 2.1, 1.4,
             "32× рост\nза 2 года\n(Латвия)",
             size=14, bold=True, color=GOLD, line_spacing=1.3)

    text_box(s, 7.0, 4.75, 5.5, 0.5,
             "Российские РЭБ (Krasukha-4, Borisoglebsk-2). Чёрное море, Восточная Европа",
             size=11, italic=True, color=MID, line_spacing=1.2)

    text_box(s, 7.0, 5.3, 5.5, 0.35, "Защита:",
             size=13, bold=True, color=DEEP)
    gps_def = ["Multi-GNSS (GPS+GLONASS+Galileo+BeiDou)",
               "INS-fallback (инерциальная)",
               "eLORAN наземная радионавигация"]
    for i, d in enumerate(gps_def):
        text_box(s, 7.15, 5.65 + i * 0.3, 5.3, 0.28, "• " + d,
                 size=11, color=MID)

    # Bottom callout
    ocean_box(s, 0.6, 6.75, 12.13, 0.45, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.8, 11.63, 0.35,
             "Spillover: военный РЭБ-эффект распространяется на не-комбатантов. Защита GNSS — collective good",
             size=11, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    add_speaker_notes(s,
        "Второй провальный кейс Sense — adversarial-атаки на SAR ATR. Идея "
        "проста: classifier обучен распознавать танки и пусковые установки по "
        "SAR-снимкам; противник размещает на технике дешёвые металлические "
        "рассеиватели в специальной геометрии, и classifier начинает "
        "неправильно классифицировать объекты. Урок. Стандартный benchmark "
        "accuracy обманчив для adversarial-доменов. Защита требует трёх вещей "
        "одновременно: Bayesian uncertainty estimates, adversarial training, и "
        "abstention pathway — модель эскалирует к человеку, а не выбирает "
        "«наиболее вероятный» класс.\n\n"
        "Третий провальный кейс Sense — GPS-spoofing гражданской авиации. AI "
        "как таковой ни при чём — но кейс демонстрирует фундаментальную "
        "хрупкость GNSS-зависимых систем. По данным Латвии, в 2024 году "
        "зарегистрировано 820 случаев интерференции против 26 в 2022 — "
        "тридцатидвухкратный рост. Атрибуция к российским средствам РЭБ — "
        "Krasukha-4, Borisoglebsk-2; зона включает Чёрное море и Восточную "
        "Европу.\n\n"
        "Урок. GNSS-only — это single point of failure. Защита — multi-GNSS, "
        "INS-fallback, eLORAN. Spillover-проблема: военный РЭБ-эффект "
        "распространяется на не-комбатантов: гражданские самолёты регулярно "
        "попадают в зоны искажённого GPS. Защита GNSS — collective good.")
    return s


def slide_15_sense_criteria(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Sense — когда не AI: 2 критерия. Distribution shift + single-sensor")

    criteria = [
        ("chart-no-axes-column", "1",
         "Low-data domain или distribution shift inevitable",
         "Если домен редкий или входное распределение меняется (новые маскировки, новые условия) — ML не выучит",
         "Альтернатива: классическая обработка сигналов + multi-sensor fusion"),
        ("unplug", "2",
         "High-stakes single-sensor decision без избыточности",
         "Если решение на одном сенсоре и стоимость ошибки высокая — AI = single point of failure",
         "Канонический контрпример: F-35 ALIS без HITL для flight authorisation"),
    ]
    cc_w = 5.95; cc_h = 4.2; gap = 0.25
    for i, (icon, num, title, desc, alt) in enumerate(criteria):
        x = 0.6 + i * (cc_w + gap)
        ocean_box(s, x, 1.85, cc_w, cc_h)
        # Badge
        ocean_box(s, x + 0.3, 2.05, 0.7, 0.7,
                  fill=GOLD, stroke=GOLD, radius_pt=20)
        text_box(s, x + 0.3, 2.05, 0.7, 0.7, "#" + num,
                 size=16, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Icon
        icon_p = ASSETS / "icons" / f"{icon}-96.png"
        if icon_p.exists():
            add_image(s, icon_p, x + cc_w - 1.2, 2.05, w=0.9, h=0.9)
        # Title
        text_box(s, x + 0.3, 3.0, cc_w - 0.6, 0.7, title,
                 size=17, bold=True, color=DEEP, line_spacing=1.2)
        # Description
        text_box(s, x + 0.3, 3.85, cc_w - 0.6, 1.1, desc,
                 size=13, color=MID, italic=True, line_spacing=1.3)
        # Alternative
        hr_line(s, x + 0.3, 5.05, cc_w - 0.6, color=LIGHT, weight=0.5)
        text_box(s, x + 0.3, 5.15, cc_w - 0.6, 0.85, "→ " + alt,
                 size=12, color=TEAL, italic=True, line_spacing=1.25, bold=True)

    # Takeaway
    ocean_box(s, 0.6, 6.35, 12.13, 0.65, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.43, 11.63, 0.5,
             "Sense — звено благополучное. Multi-sensor fusion + HITL gate — рабочая архитектура.",
             size=14, italic=True, color=DEEP, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    add_speaker_notes(s,
        "Из разобранных провалов извлекаем два критерия для звена Sense. Это "
        "первые два из общей нумерации матрицы Раздела 5.\n\n"
        "Критерий первый — распределение. Низкая плотность данных или "
        "неминуемый distribution shift. Если домен редкий или входное "
        "распределение меняется в полёте, ML-классификатор не успеет выучить. "
        "Классическая обработка сигналов плюс multi-sensor fusion стоит дешевле "
        "и даёт больше предсказуемости.\n\n"
        "Критерий второй — single-sensor. High-stakes single-sensor decision "
        "без избыточности. Если решение принимается на одном сенсоре и его "
        "стоимость ошибки высокая, AI на этом сенсоре — это single point of "
        "failure. Минимум — независимый второй канал. Канонический контрпример — "
        "ALIS без HITL для flight authorisation.\n\n"
        "Подытог раздела Sense. Самое благополучное звено OODA для AI. В "
        "коммерческой спутниковой аналитике AI уже работает в часах от съёмки "
        "до сводки; в predictive maintenance гражданской авиации AI обслуживает "
        "тысячи самолётов. Дальше переходим к самому тонкому звену цепи — "
        "Decide.")
    return s


# Section dividers and remaining slides

def slide_16(prs):
    return section_divider(prs, 2, "Decide — от наблюдения к решению",
        "Где LLM-хайп опасен, и почему «accuracy» — не та метрика",
        current_section=2,
        caption="14 минут · 5 кейсов · 3 провала (Lavender / Lancet / Vincennes) · 2 критерия")


def slide_17_decide_intro(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Decide — звено, где LLM-хайп опаснее всего", size=26)

    text_box(s, 0.6, 1.4, 12.13, 0.5,
             "«Accuracy 90%» звучит хорошо до момента, когда 10% — это тысячи человек",
             size=18, italic=True, color=MID, line_spacing=1.3)

    # Pipeline visual
    text_box(s, 0.6, 2.4, 12.13, 0.4, "Decide stack:",
             size=14, bold=True, color=DEEP)

    sources = [("Text reports", "newspaper"), ("Image", "scan"),
               ("Map", "map"), ("Telemetry", "activity")]
    sx = 0.6; sw = 2.0; gap = 0.15
    for i, (label, icon) in enumerate(sources):
        x = sx + i * (sw + gap)
        ocean_box(s, x, 2.9, sw, 1.05, fill=LIGHT_TINT)
        icon_p = ASSETS / "icons" / f"{icon}-48.png"
        if icon_p.exists():
            add_image(s, icon_p, x + 0.15, 3.05, w=0.45, h=0.45)
        text_box(s, x + 0.7, 3.1, sw - 0.8, 0.7, label,
                 size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    # arrow to fusion
    add_arrow(s, 9.3, 3.25, 0.3, 0.35, fill=MID)

    ocean_box(s, 9.7, 2.9, 1.5, 1.05, fill=MID)
    text_box(s, 9.7, 2.9, 1.5, 1.05, "Fusion",
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_arrow(s, 11.3, 3.25, 0.3, 0.35, fill=MID)

    ocean_box(s, 11.7, 2.9, 1.1, 1.05, fill=GOLD)
    text_box(s, 11.7, 2.9, 1.1, 1.05, "COA",
             size=14, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Gold callout - cost asymmetry — single wide textbox for the equation
    ocean_box(s, 0.6, 4.5, 12.13, 2.2, fill=GOLD_TINT, stroke=GOLD)
    text_runs(s, 0.85, 4.7, 11.63, 1.3, [
        {"text": "10% × 37 000 = ", "size": 60, "bold": True, "color": DEEP},
        {"text": "3 700", "size": 60, "bold": True, "color": GOLD},
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.85, 6.0, 12.0, 0.6,
             "Lavender, Газа 2023-24. 90% accuracy в life-and-death = "
             "3 700 человек, помеченных по ошибке. Метрика была не той.",
             size=13, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    text_box(s, 0.6, 6.85, 12.13, 0.3,
             "Если потенциальная цена ошибки — человеческая жизнь, «accuracy %» — "
             "это показатель количества кошмаров, которые вы готовы принять",
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Decide в OODA — это переход от наблюдения к выбору действия. Семейство "
        "задач: mission planning, target identification, multi-source fusion, "
        "decision support для командира.\n\n"
        "В отличие от Sense, где ML опирается на сенсорные features, в Decide "
        "LLM и foundation models обрабатывают смешанный input: текстовые отчёты, "
        "изображения, тактические карты, бортовые телеметрии. Это территория "
        "multi-modal foundation models — Anthropic, OpenAI, DeepSeek, Qwen, "
        "Palantir, Scale, Helsing.\n\n"
        "Главный для нашего слуха момент я хочу проиллюстрировать одним числом: "
        "10 процентов от 37 тысяч — это 3 700. Эта арифметика — это "
        "пред-просмотр канонического разбора Lavender. Если модель ошибается "
        "в 10 процентах случаев и применяется к 37 тысячам человек, ошибка "
        "масштабируется в тысячи людей. «Accuracy» как метрика проектировалась "
        "под симметрию FP и FN. В life-and-death это никогда не правда. FP — "
        "это жизнь невинного человека, и эта жизнь дороже, чем «упустить "
        "оперативника». Метрика «accuracy» обнуляется в этом контексте.")
    return s


def slide_18_palantir(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Palantir MSS — главный американский decision-support флагман")

    # Left: timeline
    ocean_box(s, 0.6, 1.95, 7.5, 4.7)
    text_box(s, 0.8, 2.1, 7.0, 0.45, "Контрактные milestones",
             size=15, bold=True, color=DEEP)
    milestones = [
        ("Май 2024", "$480M IDIQ Army contract", MID),
        ("Сентябрь 2024", "+$99,8M · расширение все рода", LIGHT),
        ("Май 2025", "+$795M ceiling increase", GOLD),
        ("До 2029", "~$1,3 миллиарда потолок", GOLD),
    ]
    for i, (date, desc, col) in enumerate(milestones):
        y = 2.7 + i * 0.85
        # Dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.95), Inches(y + 0.15),
                                 Inches(0.25), Inches(0.25))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background()
        text_box(s, 1.4, y, 2.5, 0.4, date,
                 size=14, bold=True, color=col)
        text_box(s, 4.0, y, 4.0, 0.4, desc,
                 size=13, color=DEEP, italic=True)

    # Capability bullets
    text_box(s, 0.8, 6.0, 7.0, 0.4, "Capability:",
             size=13, bold=True, color=DEEP)
    caps = ["Fusion multi-источниковой разведки",
            "AI-assisted target nomination",
            "Дашборды для командиров"]
    for i, c in enumerate(caps):
        text_box(s, 0.95, 6.35 + i * 0.0, 7.0, 0.3, " · ".join(caps),
                 size=11, italic=True, color=MID)
        break  # one line

    # Right: L1 Assistive badge
    ocean_box(s, 8.4, 1.95, 4.4, 4.7, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 8.4, 2.2, 4.4, 1.0, "L1",
             size=72, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
    text_box(s, 8.4, 3.4, 4.4, 0.5, "Assistive",
             size=22, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    hr_line(s, 8.7, 4.05, 3.8, color=GOLD, weight=1.0)
    text_box(s, 8.7, 4.25, 3.8, 1.5,
             "AI выдаёт детекции и сводки.\nКомандир решает.",
             size=15, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1.4)
    text_box(s, 8.7, 5.7, 3.8, 0.8,
             "Никаких kinetic engagements MSS не делает — решение за человеком",
             size=11, italic=True, color=MID,
             align=PP_ALIGN.CENTER, line_spacing=1.3)

    # History footer
    text_box(s, 0.6, 6.85, 12.13, 0.3,
             "История Maven 2017: Google leak март 2018 → 4000+ подписей → "
             "контракт не продлён июнь 2018 → подхвачен Palantir, Anduril, Scale (см. Раздел 4)",
             size=11, italic=True, color=LIGHT)

    add_footer(s, "Источники: DefenseScoop 2024-2025 · GovConWire 2024")
    add_speaker_notes(s,
        "Главный американский decision-support флагман — Palantir Maven Smart "
        "System. История начинается с Project Maven в 2017 году — программа "
        "DoD по анализу drone footage. В марте 2018 года через утечку стало "
        "известно, что Google помогает; к июню 2018 контракт с Google не был "
        "продлён под давлением сотрудников. Программа была подхвачена Anduril, "
        "Palantir и Scale.\n\n"
        "MSS — UI-orchestration layer Palantir над Maven AI. Контракты: "
        "первый IDIQ на 480 миллионов в мае 2024 года; дополнение на 99,8 "
        "миллиона в сентябре 2024 года; увеличение потолка на 795 миллионов "
        "в мае 2025 года. Суммарный потолок — около 1,3 миллиарда долларов до "
        "2029 года.\n\n"
        "Capability — fusion мульти-источниковой разведки, AI-assisted target "
        "nomination, дашборды. Уровень автономии — L1, Assistive. AI выдаёт "
        "детекции и сводки, командир решает. Никаких kinetic engagements MSS "
        "сам не делает.\n\n"
        "Что инженеру важно: Palantir выиграл рынок не только моделями, но "
        "инфраструктурным стеком — FedRAMP HIGH, авторизация на нескольких "
        "уровнях classified networks. Когда оцениваете defense-AI вендора, "
        "разделяйте две оси: AI capability и authorization stack.")
    return s


def slide_19_scale_helsing(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Scale AI эволюция + Helsing Altra — другие ключевые игроки Decide")

    # Left: Scale AI evolution
    ocean_box(s, 0.6, 1.95, 7.0, 4.7)
    text_box(s, 0.8, 2.1, 6.6, 0.45, "Scale AI · evolution",
             size=16, bold=True, color=DEEP)

    products = [
        ("Donovan", "2022-23", "XVIII Airborne Corps · 100k+ страниц · classified", MID),
        ("Defense Llama", "Nov 2024", "Fine-tuned Llama 3 · ops planning", LIGHT),
        ("Thunderforge", "Mar 2025", "CENTCOM + INDOPACOM · COA wargaming", GOLD),
    ]
    for i, (name, date, desc, col) in enumerate(products):
        y = 2.7 + i * 1.0
        ocean_box(s, 0.8, y, 6.6, 0.85, fill=LIGHT_TINT if col != GOLD else GOLD_TINT)
        text_box(s, 1.0, y + 0.1, 1.85, 0.4, name,
                 size=14, bold=True, color=col)
        text_box(s, 2.95, y + 0.1, 1.4, 0.4, date,
                 size=12, italic=True, color=MID)
        text_box(s, 1.0, y + 0.45, 5.5, 0.35, desc,
                 size=11, color=DEEP, italic=True)

    text_box(s, 0.8, 5.85, 6.6, 0.4,
             "Авторизация: FedRAMP HIGH · SC2S · SIPR · DISA IL4 · JWICS",
             size=10, italic=True, color=LIGHT)

    # Right: Helsing
    ocean_box(s, 7.9, 1.95, 4.9, 4.7, fill=LIGHT_TINT)
    text_box(s, 8.1, 2.1, 4.5, 0.4, "Helsing (Europe)",
             size=16, bold=True, color=DEEP)
    text_box(s, 8.1, 2.5, 4.5, 0.3, "Главный европейский игрок",
             size=11, italic=True, color=MID)

    # Big valuation number
    text_runs(s, 8.1, 3.0, 4.5, 1.1, [
        {"text": "€12 млрд", "size": 36, "bold": True, "color": GOLD},
        {"newpara": True, "text": "оценка после Series D (июнь 2025)",
         "size": 11, "italic": True, "color": MID},
    ])

    # Sub-products
    hr_line(s, 8.1, 4.25, 4.5, color=LIGHT, weight=0.5)
    helsing = [
        ("Altra", "ISR fusion + spotters для land combat targeting"),
        ("Centaur", "AI-пилот · Saab Gripen E test June 2025"),
    ]
    for i, (name, desc) in enumerate(helsing):
        y = 4.4 + i * 0.85
        text_box(s, 8.1, y, 4.5, 0.35, name,
                 size=14, bold=True, color=DEEP)
        text_box(s, 8.1, y + 0.4, 4.5, 0.5, desc,
                 size=11, color=MID, italic=True, line_spacing=1.2)

    text_box(s, 8.1, 6.2, 4.5, 0.35,
             "Главный инвестор — Prima Materia / Daniel Ek (Spotify)",
             size=10, italic=True, color=LIGHT)

    add_footer(s,
        "Источники: BusinessWire 2023 · DefenseScoop 2024 · Helsing IR 2025 · CNBC 2025")
    add_speaker_notes(s,
        "Помимо Palantir в decision-support сегменте 2026 года ещё два крупных "
        "игрока. Первый — Scale AI. Три продукта в одной линии — это эволюция. "
        "Donovan — decision-support LLM на classified networks, развёрнутый в "
        "XVIII Airborne Corps в 2022-2023 годах. 100 тысяч плюс страниц orders, "
        "SitReps, intel прошли через него. Donovan был первым LLM, развёрнутым "
        "в US classified network. В ноябре 2024 — Defense Llama, fine-tuned "
        "Llama 3 на defense corpus. И в марте 2025 — Thunderforge: wargaming "
        "плюс COA generation, используется в CENTCOM и INDOPACOM. Авторизация — "
        "FedRAMP HIGH, развёртывание на трёх уровнях classified networks.\n\n"
        "Второй — Helsing. Главный европейский игрок в defense-AI. Altra — "
        "fusion ISR-дронов плюс наблюдателей-spotters для land combat. Параллельно "
        "Helsing развивает Centaur — AI-пилота; в июне 2025 года успешно прошли "
        "испытания на Saab Gripen E. Финансы — 600 миллионов евро раунд D в "
        "июне 2025, оценка — 12 миллиардов евро. Главный инвестор — Prima "
        "Materia под управлением Дэниэла Эка, со-основателя Spotify.\n\n"
        "Если ваш карьерный интерес — европейская defense-AI индустрия, "
        "Helsing и UK Ministry of Defence — главные адреса.")
    return s


def slide_20_anthropic_russian(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Anthropic IL6 (ноябрь 2024) + Russian C2 (Svod/Glaz-Groza) — две разные карты")

    # Left: Anthropic-Palantir-AWS
    ocean_box(s, 0.6, 1.95, 6.0, 4.7, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.8, 2.1, 5.6, 0.4, "Anthropic + Palantir + AWS",
             size=16, bold=True, color=DEEP)
    text_box(s, 0.8, 2.55, 5.6, 0.35, "7 ноября 2024",
             size=12, italic=True, color=GOLD, bold=True)

    text_runs(s, 0.8, 3.05, 5.6, 0.8, [
        {"text": "Claude 3 + 3.5", "size": 18, "bold": True, "color": DEEP},
        {"text": " на ", "size": 14, "color": MID},
        {"text": "IL6", "size": 24, "bold": True, "color": GOLD},
    ])
    text_box(s, 0.8, 3.85, 5.6, 0.35,
             "Высший US gov-cloud уровень секретности",
             size=12, italic=True, color=MID)

    hr_line(s, 0.8, 4.4, 5.4, color=GOLD, weight=0.75)
    use_cases = ["Complex data processing",
                 "Pattern identification",
                 "Time-sensitive decisions"]
    for i, uc in enumerate(use_cases):
        text_box(s, 0.95, 4.55 + i * 0.35, 5.4, 0.3, "• " + uc,
                 size=12, color=DEEP, line_spacing=1.2)

    text_box(s, 0.8, 5.95, 5.6, 0.55,
             "Точка инфлексии в industry posture (см. Раздел 4.5)",
             size=11, italic=True, color=DARK_GREY, line_spacing=1.3, bold=True)

    # Right: Russian C2
    ocean_box(s, 6.9, 1.95, 5.9, 4.7, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 7.1, 2.1, 5.5, 0.4, "Russian C2: Svod / Glaz-Groza",
             size=16, bold=True, color=DEEP)
    text_box(s, 7.1, 2.55, 5.5, 0.35, "Caveat: single-source attribution",
             size=12, italic=True, color=TEAL, bold=True)

    components = [
        ("Svod Tactical Situational Awareness", "Anонс авг 2025, разработка с 2024"),
        ("Glaz", "Приложения для операторов дронов (геомэппинг)"),
        ("Groza", "Fire-control + mission management"),
        ("ZOV Maps", "Геопространственная платформа"),
    ]
    for i, (name, desc) in enumerate(components):
        y = 3.1 + i * 0.7
        text_box(s, 7.2, y, 5.4, 0.3, name,
                 size=13, bold=True, color=DEEP)
        text_box(s, 7.2, y + 0.32, 5.4, 0.3, desc,
                 size=10, italic=True, color=MID, line_spacing=1.15)

    text_box(s, 7.1, 5.95, 5.5, 0.55,
             "Independent western verification отсутствует. Effectiveness — uneven (CSIS)",
             size=10, italic=True, color=TEAL, line_spacing=1.25, bold=True)

    # Bridge bottom
    text_box(s, 0.6, 6.8, 12.13, 0.3,
             "Не упоминать был бы перекос. Некритично сообщать как success — "
             "была бы пропаганда. Промежуточный путь — упомянуть с явной оговоркой источников",
             size=10, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Четвёртый кейс — Anthropic-Palantir-AWS partnership от 7 ноября 2024. "
        "Anthropic вывел свои модели Claude 3 и 3.5 на IL6 — Impact Level 6, "
        "высший US gov-cloud уровень секретности — через Palantir и AWS GovCloud. "
        "Use cases: complex data processing, pattern identification, time-"
        "sensitive decisions. Это точка инфлексии в industry posture — между "
        "Maven walkout 2018 и этим partnership всего шесть лет, и AI-индустрия "
        "прошла полный цикл. Разбор в Разделе 4.5.\n\n"
        "Пятый кейс — российский C2. По CSIS Bondar апрель 2026, Россия строит "
        "экосистему network-centric warfare. Svod Tactical Situational Awareness "
        "Complex — объявлен в августе 2025, разработка с 2024, экспериментальное "
        "развёртывание с осени 2025. Glaz — приложения для операторов дронов. "
        "Groza — fire-control и mission management. ZOV Maps — геопространственная "
        "платформа.\n\n"
        "Caveat single-source. Information поступает из двух источников: Russian "
        "official press и CSIS-аналитика. Independent western verification "
        "отсутствует. Effectiveness — uneven. Мы упоминаем эти системы потому, "
        "что они существуют как попытка. Не упоминать был бы перекос; некритично "
        "сообщать как success — была бы пропаганда. Мы выбираем промежуточный "
        "путь: упомянуть с явной оговоркой. Это лучший паттерн работы с "
        "непроверяемыми данными вообще — не «верить» и не «отвергать», а явно "
        "маркировать уровень доказательности.")
    return s


def slide_21_lavender(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "IDF Lavender — 37 000 помечены, 90% accuracy → 3 700 false positives",
        size=22)

    # Left: cascade chart
    ocean_box(s, 0.6, 1.95, 7.4, 4.7)
    chart_p = ASSETS / "charts" / "c21-lavender-funnel.png"
    if chart_p.exists():
        add_image(s, chart_p, 0.7, 2.05, w=7.2, h=4.0)
    text_box(s, 0.7, 6.15, 7.2, 0.4,
             "20 секунд review per target · 15-20 civilian casualties auth'd per junior operative",
             size=10, italic=True, color=MID, line_spacing=1.2, align=PP_ALIGN.CENTER)

    # Right: 3 lessons
    lessons = [
        ("«Accuracy» — wrong metric",
         "FP consequence × population × frequency.\nCost-asymmetry FP↔FN неприемлемо"),
        ("AI снимает фрикцию",
         "Темпы вырастают → качество deliberation падает.\nВ life-and-death катастрофично"),
        ("HITL ≠ Human-In-Decision",
         "20 sec review = формальный HITL,\nфункциональный HOTL"),
    ]
    lc_w = 4.7; lc_h = 1.45; gap = 0.12
    for i, (title, desc) in enumerate(lessons):
        y = 1.95 + i * (lc_h + gap)
        ocean_box(s, 8.1, y, lc_w, lc_h, fill=GOLD_TINT if i == 0 else SURFACE,
                  stroke=GOLD if i == 0 else LIGHT)
        text_box(s, 8.3, y + 0.15, lc_w - 0.4, 0.4,
                 f"Урок {i+1}: {title}",
                 size=13, bold=True, color=DEEP)
        text_box(s, 8.3, y + 0.6, lc_w - 0.4, 0.85, desc,
                 size=11, color=MID, italic=True, line_spacing=1.25)

    # Source footer
    add_footer(s,
        "Источники: Abraham, +972 / Local Call (апрель 2024); ICRC, Lieber Institute, AOAV — academic разборы; IDF опровергает")

    add_speaker_notes(s,
        "Главный педагогический провал звена Decide — IDF Lavender, "
        "AI-система массовой идентификации целей в Газе 2023-2024.\n\n"
        "Что произошло. Lavender — AI-database, помечающая палестинских мужчин "
        "как «подозреваемых» по паттернам коммуникации. Около 37 000 человек "
        "помечено. По собственному признанию ЦАХАЛ, точность около 90% — то "
        "есть около 3 700 человек false positive. Среднее время review одной "
        "цели — около 20 секунд. Авторизованный «сопутствующий ущерб» — до "
        "15-20 гражданских жертв на одного оперативника низшего звена.\n\n"
        "Уроки три.\n\n"
        "Урок первый. «Accuracy %» — не та метрика для life-and-death. Правильная "
        "метрика — false positive consequence × population × frequency. В "
        "медицине false negative дороже false positive — система настраивается "
        "на высокую sensitivity. В Lavender — наоборот: FP — жизнь невинного "
        "человека. Но система проектировалась под симметрию.\n\n"
        "Урок второй. AI снимает фрикцию decision-making → темпы вырастают → "
        "качество deliberation падает. Снимая фрикцию, AI масштабирует не "
        "качество решения, а его скорость. В life-and-death катастрофично.\n\n"
        "Урок третий. Human-in-the-loop не равно human-in-the-decision. 20 "
        "секунд проверки — это не review, это формальное подтверждение. Если "
        "HITL вырождается в подпись без осмысления, это HOTL под маской HITL. "
        "Engineering decision «сколько времени у оператора» — это формальная "
        "категоризация системы.\n\n"
        "Альтернатива не «давайте сделаем Lavender точнее». Альтернатива — "
        "изменение архитектуры: AI ассистирует triage, human keeps authority "
        "с реальным временем на review. AI — accelerator, не decision-maker.")
    return s


def slide_22_lancet_vincennes(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Lancet rollback (демо ≠ продакшен) + Vincennes 1988 (UI под стрессом → LLM)",
        size=22)

    # Left: Lancet
    ocean_box(s, 0.6, 1.95, 6.0, 4.7)
    text_box(s, 0.8, 2.1, 5.6, 0.4, "Russian Lancet ATR rollback",
             size=16, bold=True, color=DEEP)
    text_box(s, 0.8, 2.5, 5.6, 0.35, "2022-2024 · LO2 canonical case",
             size=11, italic=True, color=GOLD, bold=True)

    text_box(s, 0.8, 3.05, 5.6, 0.35, "2022-23:", size=12, bold=True, color=MID)
    text_box(s, 1.4, 3.05, 5.0, 0.35,
             "«autonomously find and hit target» + videos с UI",
             size=11, italic=True, color=DEEP)

    text_box(s, 0.8, 3.55, 5.6, 0.35, "2024:", size=12, bold=True, color=GOLD)
    text_box(s, 1.4, 3.55, 5.0, 0.35,
             "CSIS / MWI: AI-guidance off · videos без autonomous-locking UI",
             size=11, italic=True, color=DEEP)

    icon_p = ASSETS / "icons" / "triangle-alert-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 5.9, 2.05, w=0.55, h=0.55)

    text_box(s, 0.8, 4.2, 5.6, 0.5,
             "Edge cases — это БОЛЬШАЯ ЧАСТЬ реального поля боя: "
             "пыль, дым, EW, новые маскировки",
             size=11, color=DEEP, italic=True, line_spacing=1.25)

    hr_line(s, 0.8, 4.85, 5.4, color=LIGHT, weight=0.5)
    text_box(s, 0.8, 4.95, 5.6, 0.4, "Альтернатива:",
             size=12, bold=True, color=DEEP)
    text_box(s, 0.8, 5.3, 5.6, 0.5,
             "Operator-in-the-loop + automated tracking-assist. Не autonomous engage до production hardening",
             size=11, color=TEAL, italic=True, line_spacing=1.25)

    text_box(s, 0.8, 6.05, 5.6, 0.4,
             "ML perf в narrow distribution не переносится на full battlefield variance",
             size=11, bold=True, color=DEEP, italic=True, line_spacing=1.2)

    # Right: Vincennes
    ocean_box(s, 6.8, 1.95, 5.95, 4.7)
    text_box(s, 7.0, 2.1, 5.6, 0.4, "USS Vincennes / Iran Air 655",
             size=16, bold=True, color=DEEP)
    text_box(s, 7.0, 2.5, 5.6, 0.35, "3 июля 1988 · 290 KIA",
             size=11, italic=True, color=RED_WARN, bold=True)

    icon_p = ASSETS / "icons" / "eye-off-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 12.0, 2.05, w=0.55, h=0.55)

    events = [
        ("Aegis записал track как climbing", "✓", TEAL),
        ("Экипаж под стрессом доложил «descending into attack»", "✗", RED_WARN),
        ("2 ракеты SM-2 → 290 KIA", "", GOLD),
    ]
    for i, (txt, mark, col) in enumerate(events):
        y = 3.0 + i * 0.55
        if mark:
            text_box(s, 7.05, y, 0.4, 0.4, mark,
                     size=18, bold=True, color=col)
        text_box(s, 7.5, y, 5.1, 0.4, txt,
                 size=11, color=DEEP, italic=True, line_spacing=1.2)

    hr_line(s, 7.0, 4.75, 5.5, color=LIGHT, weight=0.5)

    # Bridge to LLM
    ocean_box(s, 7.0, 4.9, 5.5, 1.45, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 7.15, 5.0, 5.2, 0.35, "Bridge to LLM:",
             size=12, bold=True, color=DEEP)
    text_box(s, 7.15, 5.35, 5.2, 0.95,
             "LLM выдаёт fluent confident output → оператор под давлением "
             "склонен принять. Confident BS = high-risk confident BS в high-stakes",
             size=11, color=DEEP, italic=True, line_spacing=1.25)

    add_footer(s,
        "Источники: CSIS 2025 Lancet · USNI Proceedings Jul 2018 · Foreign Affairs 2024")
    add_speaker_notes(s,
        "Второй провал звена Decide — Russian Lancet ATR rollback, 2022-2024. "
        "Канонический кейс для LO2: отличить демо от продакшена.\n\n"
        "Маркетинг 2022-23: «autonomously find and hit target», видео с "
        "интерфейсом «Target Locked». Анализ CSIS / MWI 2024: AI-guidance off "
        "после первоначальных развёртываний. Последние video drops без "
        "autonomous-locking UI. Гипотеза: premature product rollout с "
        "последующим recall. ATR работал в demo-conditions, но не в реальных: "
        "пыль, дым, EW, новые маскировки. Это урок не Lancet-specific. ML "
        "performance в narrow distribution не переносится на full battlefield. "
        "Применимо ко всем ATR-системам, drone autonomy claims, autonomous "
        "targeting рекламам.\n\n"
        "Третий провал — USS Vincennes 1988. 3 июля крейсер сбил иранский "
        "гражданский Airbus A300, погибли 290 человек. Aegis-система корректно "
        "записала track как climbing — характеристика гражданского самолёта. "
        "Операторы под стрессом доложили «descending into attack». Не баг "
        "алгоритма — automation выполнил работу. Сбой на интерфейсе под combat "
        "stress.\n\n"
        "Почему это про LLM. В современных LLM-системах структурно похожая "
        "проблема: LLM выдаёт fluent confident output — текст хорошо написан "
        "и звучит уверенно, — и оператор под временным давлением склонен принять "
        "его, не верифицируя источники. Confident BS = high-risk confident BS "
        "в high-stakes. Урок 1988 года прямо применим к 2026.")
    return s


def slide_23_decide_criteria(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, "Decide — когда не AI: 2 критерия")

    criteria = [
        ("circle-help", "3",
         "Long-tail edge cases с low ML confidence",
         "Automation bias масштабирует ошибки. Нужен structured abstention — AI говорит «не знаю» и эскалирует",
         "Calibrated uncertainty + explicit threshold + UI"),
        ("shield-alert", "4",
         "High-stakes life-and-death без HITL",
         "Cost-asymmetry FP↔FN слишком велика для статистики. Lavender — канонический контрпример",
         "Real HITL — не 20-сек подпись. ms-на-review = формальная категоризация системы"),
    ]
    cc_w = 5.95; cc_h = 4.2; gap = 0.25
    for i, (icon, num, title, desc, alt) in enumerate(criteria):
        x = 0.6 + i * (cc_w + gap)
        ocean_box(s, x, 1.85, cc_w, cc_h)
        ocean_box(s, x + 0.3, 2.05, 0.7, 0.7,
                  fill=GOLD, stroke=GOLD, radius_pt=20)
        text_box(s, x + 0.3, 2.05, 0.7, 0.7, "#" + num,
                 size=16, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        icon_p = ASSETS / "icons" / f"{icon}-96.png"
        if icon_p.exists():
            add_image(s, icon_p, x + cc_w - 1.2, 2.05, w=0.9, h=0.9)
        text_box(s, x + 0.3, 3.0, cc_w - 0.6, 0.7, title,
                 size=17, bold=True, color=DEEP, line_spacing=1.2)
        text_box(s, x + 0.3, 3.85, cc_w - 0.6, 1.1, desc,
                 size=13, color=MID, italic=True, line_spacing=1.3)
        hr_line(s, x + 0.3, 5.05, cc_w - 0.6, color=LIGHT, weight=0.5)
        text_box(s, x + 0.3, 5.15, cc_w - 0.6, 0.85, "→ " + alt,
                 size=12, color=TEAL, italic=True, bold=True, line_spacing=1.25)

    ocean_box(s, 0.6, 6.35, 12.13, 0.65, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.43, 11.63, 0.5,
             "Decide — звено самое тонкое. LLM-hype опаснее всего. AI — accelerator, не decision-maker.",
             size=14, italic=True, color=DEEP, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_notes(s,
        "Из разобранных провалов извлекаем два критерия для звена Decide — "
        "критерии 3 и 4.\n\n"
        "Критерий третий — long-tail. Long-tail edge cases с низкой ML "
        "confidence. Если задача — decision в области, где модель часто "
        "встречается с примерами вне обучающего распределения, automation bias "
        "масштабирует ошибки. Здесь нужен structured abstention: AI говорит «не "
        "знаю» и эскалирует. Engineering: calibrated uncertainty, explicit "
        "threshold для эскалации, UI который показывает неуверенность.\n\n"
        "Критерий четвёртый — life-and-death без HITL. Cost-asymmetry FP↔FN "
        "слишком велика для чисто статистического решения. Lavender — "
        "канонический контрпример. Формальный HITL обязателен — и должен быть "
        "реальным, не вырожденным в 20-сек подпись.\n\n"
        "Подытог Раздела 2. Decide — звено, где LLM-хайп опаснее всего. "
        "Доступны мощные инструменты — MSS, Donovan, Defense Llama, Helsing, "
        "Anthropic IL6. Параллельно — Lavender показывает wrong metric; "
        "Lancet — demo ≠ production; Vincennes — UI под стрессом и confident "
        "BS. Два критерия — long-tail и life-and-death без real HITL.")
    return s


# Continue with rest...
def slide_24(prs):
    return section_divider(prs, 3, "Act — автономия на платформе",
        "Где hype далеко впереди реальности — но adoption всё-таки растёт",
        current_section=3,
        caption="14 минут · 6 кейсов · 3 провала (MCAS / Patriot / Replicator) · 2 критерия · RU dual-use")


def slide_25_act_intro(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Act — звено, где hype далеко впереди реальности")

    text_box(s, 0.6, 1.4, 12.13, 0.5,
             "Большая часть combat-strikes — operator-in-loop или semi-auto, не fully-autonomous",
             size=16, italic=True, color=MID)

    # L1-L5 mini-preview
    text_box(s, 0.6, 2.4, 12.13, 0.35, "Лестница автономии (preview):",
             size=14, bold=True, color=DEEP)

    levels = [
        ("L1", "Assistive", LIGHT),
        ("L2", "Semi-auto", LIGHT),
        ("L3", "Supervised", GOLD),
        ("L4", "Pre-authorised", MID),
        ("L5", "Full LAWS", RED_WARN),
    ]
    lw = 2.3; gap = 0.15
    for i, (lvl, name, col) in enumerate(levels):
        x = 0.6 + i * (lw + gap)
        is_pivot = (lvl == "L3")
        ocean_box(s, x, 2.9, lw, 1.1, fill=GOLD_TINT if is_pivot else LIGHT_TINT,
                  stroke=col)
        text_box(s, x, 2.95, lw, 0.5, lvl,
                 size=22, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        text_box(s, x, 3.45, lw, 0.5, name,
                 size=12, italic=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    text_box(s, 0.6, 4.1, 12.13, 0.35,
             "Большая часть Act в 2026 — L2-L3, не L4-L5",
             size=12, italic=True, color=MID, align=PP_ALIGN.CENTER)

    # Cost-asymmetry callout
    ocean_box(s, 0.6, 4.7, 12.13, 2.2, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.85, 4.85, 5.5, 1.4, "$300",
             size=72, bold=True, color=TEAL,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 6.4, 4.95, 0.6, 1.2, "↔",
             size=56, bold=True, color=MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 7.1, 4.85, 2.7, 1.4, "$3M",
             size=72, bold=True, color=GOLD,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.85, 6.15, 12.0, 0.6,
             "Counter-drone asymmetry — дешёвый дрон против Patriot. "
             "Замена $3M на AI-perimeter-defence — explosive growth Act 2026",
             size=13, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    add_speaker_notes(s,
        "Act в OODA — это последнее звено: AI исполняет действие, не только "
        "наблюдает или рекомендует. Pilot-class autonomy, swarm coordination, "
        "counter-drone autonomy, industrial robotics.\n\n"
        "Adoption растёт быстро по числу платформ, но большая часть combat-"
        "strikes остаётся operator-in-loop или semi-auto terminal guidance. "
        "«AI заменит пилотов» — overhyped; «AI заменит наводчиков» — "
        "overhyped; «AI заменит командиров» — overhyped. Реальность — "
        "collaborative pilots с supervised AI.\n\n"
        "Особое место в Act — counter-drone autonomy. Главная причина роста — "
        "асимметрия: 300-долларовый дрон против 3-миллионной зенитной ракеты "
        "Patriot. Заменить три миллиона на сравнимый по цене AI-perimeter-"
        "defence — это инженерная и экономическая необходимость, и это "
        "пространство explosive growth в 2026.\n\n"
        "Когда вы оцениваете defense-AI вакансию или продукт, держите в голове "
        "эту cost-asymmetry. Если экономика боя сместилась в сторону «дешёвый "
        "attacker против дорогой defence», AI оказывается основным механизмом "
        "перестройки.")
    return s


def slide_26_fury(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Anduril Fury YFQ-44A — CCA Increment 1, L3 Supervised", size=24)

    # Left: Fury silhouette + spec
    ocean_box(s, 0.6, 1.95, 7.5, 4.7)
    # Drone-aircraft outline using shapes
    text_box(s, 0.8, 2.1, 7.1, 0.4, "YFQ-44A · Fury",
             size=18, bold=True, color=DEEP)
    text_box(s, 0.8, 2.5, 7.1, 0.35,
             "Anduril CCA Increment 1 · first flight 31 Oct 2025",
             size=12, italic=True, color=MID)

    # Aircraft silhouette built via shapes
    # Body
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(3.4), Inches(5.0), Inches(0.7))
    body.adjustments[0] = 0.4
    body.fill.solid(); body.fill.fore_color.rgb = MID
    body.line.fill.background()
    # Nose
    nose = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                              Inches(6.4), Inches(3.45), Inches(0.6), Inches(0.6))
    nose.fill.solid(); nose.fill.fore_color.rgb = MID
    nose.line.fill.background()
    # Wings (delta)
    wing1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                               Inches(2.5), Inches(3.05), Inches(2.0), Inches(0.4))
    wing1.fill.solid(); wing1.fill.fore_color.rgb = LIGHT
    wing1.line.fill.background()
    wing2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                               Inches(2.5), Inches(4.05), Inches(2.0), Inches(0.4))
    wing2.rotation = 180
    wing2.fill.solid(); wing2.fill.fore_color.rgb = LIGHT
    wing2.line.fill.background()
    # Tail
    tail = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                              Inches(1.2), Inches(3.0), Inches(0.5), Inches(0.5))
    tail.rotation = 90
    tail.fill.solid(); tail.fill.fore_color.rgb = MID
    tail.line.fill.background()

    text_box(s, 0.8, 5.1, 7.0, 0.4, "Спецификации:",
             size=13, bold=True, color=DEEP)
    specs = [
        "Высота: до 50 000 футов",
        "Скорость: M 0.95 · Перегрузка 9g",
        "Двигатель: Williams FJ44-4M (4 000 фунтов тяги)",
        "Вооружение: AIM-120 AMRAAM",
    ]
    for i, sp in enumerate(specs):
        text_box(s, 0.95, 5.5 + i * 0.3, 6.8, 0.27, "• " + sp,
                 size=11, color=MID, line_spacing=1.15)

    # Right: L3 + production
    ocean_box(s, 8.4, 1.95, 4.4, 2.2, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 8.4, 2.05, 4.4, 1.0, "L3",
             size=72, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
    text_box(s, 8.4, 3.15, 4.4, 0.4, "Supervised autonomy",
             size=15, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER)
    text_box(s, 8.6, 3.6, 4.0, 0.55,
             "AI executes в pre-authorised envelope; пилотируемый wingman supervises",
             size=10, italic=True, color=MID,
             align=PP_ALIGN.CENTER, line_spacing=1.3)

    ocean_box(s, 8.4, 4.3, 4.4, 1.05, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 8.55, 4.4, 4.1, 0.4, "23 марта 2026",
             size=14, bold=True, color=DEEP)
    text_box(s, 8.55, 4.75, 4.1, 0.5,
             "Arsenal-1, Огайо · $1 млрд инвестиций",
             size=11, italic=True, color=MID, line_spacing=1.2)

    ocean_box(s, 8.4, 5.5, 4.4, 1.15, fill=LIGHT_TINT)
    text_box(s, 8.55, 5.6, 4.1, 0.4, "Autonomy stack",
             size=13, bold=True, color=DEEP)
    text_box(s, 8.55, 6.0, 4.1, 0.6,
             "Hivemind (Shield AI) + Lattice (Anduril OS)",
             size=11, color=MID, italic=True, line_spacing=1.2)

    add_footer(s,
        "Источники: Wikipedia YFQ-44 · Air & Space Forces 2026 · The Aviationist March 24 2026")
    add_speaker_notes(s,
        "Anduril Fury YFQ-44A — Beautiful poster child современной американской "
        "ставки на автономию. CCA — Collaborative Combat Aircraft — программа "
        "ВВС США по созданию беспилотных wingmen, летающих рядом с пилотируемыми "
        "истребителями. Fury YFQ-44A — Increment 1.\n\n"
        "Спецификации: высота до 50 000 футов, скорость 0,95 Маха, перегрузка "
        "9g, двигатель Williams FJ44-4M на 4 000 фунтов тяги. Главное вооружение — "
        "AIM-120 AMRAAM. Первый полёт состоялся 31 октября 2025. Серийное "
        "производство стартовало 23 марта 2026 на новом заводе Arsenal-1 в "
        "Огайо — миллиард долларов инвестиций.\n\n"
        "Управляется Fury стеком: Hivemind (Shield AI autonomy stack) плюс "
        "Lattice (Anduril proprietary OS для autonomous mesh-coordination). "
        "Уровень автономии — L3, Supervised autonomy: AI executes в pre-"
        "authorised envelope, wingman сверху supervises.\n\n"
        "Что вынести: 2026 год — первый год CCA serial production. Стандартом "
        "для defense-AI в воздушных платформах становится L3 Supervised: AI "
        "делает работу, человек сверху смотрит, человек авторизует key decisions.")
    return s


# Helper for remaining slides — keeping each compact
def slide_27_x62a_saker(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "X-62A VISTA (Sep 2023 dogfight) + Saker Scout (Ukraine combat-tested)",
        size=22)

    # Left: X-62A
    ocean_box(s, 0.6, 1.95, 6.0, 4.7)
    text_box(s, 0.8, 2.1, 5.6, 0.4, "DARPA X-62A VISTA",
             size=16, bold=True, color=DEEP)
    text_box(s, 0.8, 2.5, 5.6, 0.3,
             "Модифицированный F-16 · AI-agent управляет",
             size=11, italic=True, color=MID)

    events = [
        ("Dec 2022", "Начало испытаний", MID),
        ("Feb 2023", "12 полётов в Эдвардсе", MID),
        ("Sep 2023", "Первый AI-vs-manned dogfight · 2 000 ft @ 1 200 mph", GOLD),
        ("May 2024", "USAF Secretary Kendall летал", MID),
    ]
    for i, (date, ev, col) in enumerate(events):
        y = 3.0 + i * 0.42
        text_box(s, 0.95, y, 1.5, 0.3, date,
                 size=11, bold=True, color=col)
        text_box(s, 2.45, y, 4.0, 0.3, ev,
                 size=11, color=DEEP, italic=True)

    text_box(s, 0.8, 4.85, 5.6, 0.35, "Объём: 100 000+ строк FCS changes · 21 полёт",
             size=11, italic=True, color=LIGHT)

    # Anti-hype caveat in gold-tint
    ocean_box(s, 0.8, 5.3, 5.55, 1.25, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.95, 5.4, 5.3, 0.35, "Anti-hype:",
             size=12, bold=True, color=DEEP)
    text_box(s, 0.95, 5.75, 5.3, 0.8,
             "Narrow scripted scenario · 1-на-1 dogfight · BVR исключён · "
             "fuel mgmt не покрыт · ROE не учитывался",
             size=10, italic=True, color=DEEP, line_spacing=1.25)

    # Right: Saker
    ocean_box(s, 6.8, 1.95, 5.95, 4.7)
    text_box(s, 7.0, 2.1, 5.5, 0.4, "Saker Scout (Украина)",
             size=16, bold=True, color=DEEP)
    text_box(s, 7.0, 2.5, 5.5, 0.3,
             "Combat-tested · L2 Semi-auto",
             size=11, italic=True, color=GOLD, bold=True)

    saker_stats = [
        ("64", "автономных целей идентифицирует"),
        ("~10 км", "дальность"),
        ("EW-resistant", "CV-классификация target ID"),
        ("Brave1", "300+ AI dev · 70+ AI/CV в combat"),
    ]
    for i, (big, desc) in enumerate(saker_stats):
        y = 3.0 + i * 0.55
        text_box(s, 7.05, y, 1.5, 0.4, big,
                 size=15, bold=True, color=GOLD if i == 0 else DEEP)
        text_box(s, 8.6, y + 0.05, 4.0, 0.35, desc,
                 size=11, color=MID, italic=True)

    hr_line(s, 7.0, 5.3, 5.5, color=LIGHT, weight=0.5)
    text_box(s, 7.0, 5.4, 5.5, 0.35, "Пивот 2024-25:",
             size=12, bold=True, color=DEEP)
    pivots = [
        "Dec 2024: первый unmanned ground op (UGV + FPV)",
        "2025: AI-mother-drone · 2 AI-FPV strike · 300 км",
    ]
    for i, p in enumerate(pivots):
        text_box(s, 7.15, 5.75 + i * 0.3, 5.3, 0.28, "• " + p,
                 size=10, color=MID, italic=True)

    add_footer(s,
        "Источники: DARPA 2024 · The Aviationist 2024 · MWI 2025 · CSIS 2025 · Kyiv Independent 2025")
    add_speaker_notes(s,
        "DARPA ACE X-62A VISTA — первый в мире AI-vs-человек воздушный бой в "
        "реальном времени. Хронология: декабрь 2022 — начало испытаний; "
        "февраль 2023 — 12 полётов в Эдвардсе; сентябрь 2023 — первый AI-vs-"
        "manned F-16 dogfight: 2 000 футов nose-to-nose на 1 200 миль/ч; май "
        "2024 — секретарь USAF Кендалл лично летал в AI-управляемом X-62A. "
        "100 000+ строк flight-critical software; 21 испытательный полёт.\n\n"
        "Anti-hype: X-62A — narrow scripted scenario. Один-на-один dogfight в "
        "известной зоне, BVR-критерии исключены, fuel management не покрыт, "
        "ROE не учитывался. «AI заменит пилотов» в 2026 — маркетинговая "
        "экстраполяция от узкого demo. Реальность ближе к collaborative "
        "supervised CCA (Fury) — wingman, не replacement.\n\n"
        "Saker Scout — один из combat-tested AI-loitering munitions. "
        "Идентифицирует до 64 целей; дальность около 10 км; CV-классификация "
        "для target ID; EW-resistant. Часть Brave1 — государственной платформы "
        "Украины (300+ AI dev). Первый полностью unmanned ground operation — "
        "декабрь 2024 (UGV + FPV без пехоты). В 2025 — AI-mother-drone, "
        "доставляющая 2 AI-FPV strike дрона за 300 км.\n\n"
        "Уровень — L2, Semi-auto perception. AI рекомендует target lock, "
        "оператор подтверждает. Именно Saker, а не X-62A, представляет настоящее "
        "лицо Act-автономии в 2026 году: L2, не L4-L5.")
    return s


def slide_28_geran_cognitive(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Russian Act — Geran-2 (defense) + Cognitive Pilot (civilian dual-use)",
        size=22)

    # Left: Geran-2
    ocean_box(s, 0.6, 1.95, 6.0, 4.7)
    text_box(s, 0.8, 2.1, 5.6, 0.4, "Geran-2 evolution",
             size=16, bold=True, color=DEEP)

    geran_stats = [
        ("Алабуга ОЭЗ", "Производство в Татарстане"),
        ("~2 700-3 000 / мес", "К концу 2025 [VFY-day-of]"),
        (">26 000", "произведено к поздней весне 2025"),
        (">40 000", "план к концу 2025 [VFY]"),
    ]
    for i, (big, desc) in enumerate(geran_stats):
        y = 2.6 + i * 0.45
        text_box(s, 1.0, y, 2.5, 0.35, big,
                 size=13, bold=True, color=GOLD if i in [1, 2, 3] else DEEP)
        text_box(s, 3.6, y + 0.03, 2.8, 0.3, desc,
                 size=10, italic=True, color=MID)

    hr_line(s, 0.8, 4.5, 5.4, color=LIGHT, weight=0.5)
    text_box(s, 0.8, 4.6, 5.4, 0.35, "AI-stack (wreckage analysis):",
             size=12, bold=True, color=DEEP)
    stack = ["NVIDIA Jetson onboard",
             "High-res камеры + thermal",
             "FPGA для EW-resistance",
             "2026 — anti-radiation seeker"]
    for i, s_ in enumerate(stack):
        text_box(s, 0.95, 4.95 + i * 0.25, 5.2, 0.23, "• " + s_,
                 size=10, color=MID, line_spacing=1.15)

    # Supply chain warning
    ocean_box(s, 0.8, 6.0, 5.55, 0.7, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.95, 6.1, 5.25, 0.55,
             "Supply-chain caveat: 1 111 Dell PowerEdge XE9680 через "
             "Shreya Life Sciences (India) → Russia, Apr-Aug 2024",
             size=10, italic=True, color=DEEP, line_spacing=1.2)

    # Right: Cognitive Pilot
    ocean_box(s, 6.8, 1.95, 5.95, 4.7, fill=LIGHT_TINT)
    text_box(s, 7.0, 2.1, 5.5, 0.4, "Cognitive Pilot · civilian dual-use",
             size=16, bold=True, color=DEEP)
    text_box(s, 7.0, 2.5, 5.5, 0.3,
             "JV Сбер + Cognitive Technologies (Москва)",
             size=11, italic=True, color=MID)

    text_box(s, 7.0, 3.0, 5.5, 0.35, "Stack:",
             size=13, bold=True, color=DEEP)
    text_box(s, 7.15, 3.35, 5.3, 0.4,
             "CV + радар + LiDAR (автономия без GNSS)",
             size=11, color=MID, italic=True)

    text_box(s, 7.0, 3.85, 5.5, 0.35, "Применения:",
             size=13, bold=True, color=DEEP)
    apps = ["КАМАЗ-комбайны, тракторы СберАгро",
            "Городской транспорт",
            "Железная дорога",
            "Снегоуборочная техника"]
    for i, a in enumerate(apps):
        text_box(s, 7.15, 4.2 + i * 0.28, 5.3, 0.25, "• " + a,
                 size=11, color=MID, line_spacing=1.15)

    text_runs(s, 7.0, 5.5, 5.5, 0.55, [
        {"text": "До 50 000", "size": 22, "bold": True, "color": GOLD},
        {"text": " систем/год (план)", "size": 13, "color": DEEP},
    ])

    text_box(s, 7.0, 6.15, 5.5, 0.45,
             "НЕ identified as defense supplier в открытых источниках. "
             "Те же CV+LiDAR — в CAD/CAM (Лек 6) и пром автоматизации (Лек 14)",
             size=10, italic=True, color=TEAL, line_spacing=1.25)

    add_footer(s,
        "Источники: CSIS Bondar 2026 · Tom's Hardware 2024 · Fortune 2026 · TASS 2024 · Cognitive Pilot 2025")
    add_speaker_notes(s,
        "Geran-2 — российская модификация иранского Shahed-136, производится на "
        "Алабугской ОЭЗ. К концу 2025 — около 2 700-3 000 в месяц. Общий объём >"
        "26 000 к поздней весне 2025, план >40 000 к концу года. Эти цифры — "
        "[VFY-day-of].\n\n"
        "AI-эволюция: wreckage-анализ показывает NVIDIA Jetson onboard, "
        "high-res камеры, тепловизионные модули, FPGA для EW-resistance. В 2026 "
        "появился anti-radiation seeker.\n\n"
        "Caveat по «автономии». Wreckage подтверждает onboard ML, но реальная "
        "роль autonomous decision quality vs operator override — unclear. "
        "Большая часть strikes — operator-guided + GPS-guided. «Autonomy» в "
        "смысле «решает сама» — overstated.\n\n"
        "Supply-chain каверз. 1 111 серверов Dell PowerEdge XE9680 через "
        "индийскую Shreya Life Sciences в апреле-августе 2024. Это и риск для "
        "адверсаров, и инженерный урок: hardware supply-chain — strategic risk.\n\n"
        "Cognitive Pilot — civilian dual-use balance. JV Сбер плюс Cognitive "
        "Technologies. Применения: КАМАЗ-комбайны, тракторы СберАгро, городской "
        "транспорт, железная дорога. Stack: CV + радар + LiDAR. План — до 50 000 "
        "систем/год. НЕ identified as defense supplier. Те же стеки работают в "
        "CAD/CAM (Лекция 6) и в промышленной автоматизации (Лекция 14).")
    return s


def slide_29_mcas(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Boeing 737 MAX MCAS — canonical anti-pattern safety-critical AI", size=22)

    # Left: 4 lessons in 2x2
    lessons = [
        ("Single-point-of-failure",
         "Одна модель, один сенсор, одно решение"),
        ("Opacity",
         "Пилоты НЕ знали о MCAS · нет override"),
        ("Software cures hardware",
         "MCAS компенсировал hardware shortfall"),
        ("FMEA / FTA не пройден",
         "SPOF должен быть пойман на анализе"),
    ]
    lc_w = 3.5; lc_h = 1.9; gap = 0.15
    for i, (title, desc) in enumerate(lessons):
        row = i // 2; col = i % 2
        x = 0.6 + col * (lc_w + gap)
        y = 1.95 + row * (lc_h + gap)
        ocean_box(s, x, y, lc_w, lc_h, fill=GOLD_TINT, stroke=GOLD)
        text_box(s, x + 0.15, y + 0.15, 0.5, 0.5, "✗",
                 size=24, bold=True, color=RED_WARN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + 0.65, y + 0.2, lc_w - 0.75, 0.45, title,
                 size=13, bold=True, color=DEEP, line_spacing=1.2)
        text_box(s, x + 0.65, y + 0.7, lc_w - 0.75, 1.1, desc,
                 size=11, color=MID, italic=True, line_spacing=1.25)

    # Right: vertical timeline
    ocean_box(s, 7.8, 1.95, 5.0, 4.0)
    text_box(s, 8.0, 2.05, 4.8, 0.4, "Crash timeline",
             size=15, bold=True, color=DEEP)

    crash_events = [
        ("29 окт 2018", "Lion Air 610", "189 KIA"),
        ("10 мар 2019", "Ethiopian Airlines 302", "157 KIA"),
        ("", "→ 346 KIA total", "20 мес остановки в США"),
    ]
    for i, (date, ev, casualty) in enumerate(crash_events):
        y = 2.55 + i * 0.9
        if date:
            text_box(s, 8.0, y, 1.8, 0.35, date,
                     size=12, bold=True, color=MID)
            text_box(s, 8.0, y + 0.35, 4.6, 0.3, ev,
                     size=11, color=DEEP, italic=True)
            text_box(s, 8.0, y + 0.65, 4.6, 0.3, "→ " + casualty,
                     size=10, italic=True, color=RED_WARN)
        else:
            ocean_box(s, 8.0, y, 4.6, 0.85, fill=GOLD_TINT, stroke=GOLD)
            text_box(s, 8.0, y, 4.6, 0.45, ev,
                     size=18, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            text_box(s, 8.0, y + 0.45, 4.6, 0.35, casualty,
                     size=10, italic=True, color=DEEP,
                     align=PP_ALIGN.CENTER)

    # Patriot mini-callback
    ocean_box(s, 7.8, 6.05, 5.0, 1.0, fill=TEAL_TINT, stroke=TEAL)
    icon_p = ASSETS / "icons" / "radar-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 7.9, 6.15, w=0.5, h=0.5)
    text_box(s, 8.5, 6.1, 4.2, 0.35, "Patriot mini-callback",
             size=12, bold=True, color=DEEP)
    text_box(s, 8.5, 6.4, 4.2, 0.6,
             "2003 (RAF Tornado + USN F/A-18) + 2024 (Ukr F-16). Automation bias.",
             size=10, italic=True, color=MID, line_spacing=1.25)

    add_footer(s,
        "Источники: PMC 2020 · ThinkReliability 2019 · Trenchart 2018 · SOFREP 2003")
    add_speaker_notes(s,
        "Первый и канонический провал Act — Boeing 737 MAX MCAS. Два крушения: "
        "Lion Air 610 (29 октября 2018, 189 KIA) и Ethiopian Airlines 302 (10 "
        "марта 2019, 157 KIA). Суммарно 346 погибших. 20-месячная остановка "
        "эксплуатации в США; международная un-grounding до 2022.\n\n"
        "Что произошло. Boeing 737 MAX получил большие двигатели, что сместило "
        "аэродинамический центр. Решение Boeing — программное: MCAS, "
        "автоматически корректирующая trim вниз. MCAS активировался по одному "
        "AoA-сенсору, без резервирования. Когда сенсор давал ложное показание, "
        "пилот не мог override — не было ни тренировок, ни понимания, что "
        "система делает.\n\n"
        "Строго говоря, MCAS — не AI: это classical control с if-then логикой. "
        "Но pedagogically — canonical anti-pattern для всех safety-critical AI.\n\n"
        "Уроки. Один: redundancy — никогда не делать safety-critical системы "
        "зависимыми от single sensor. Второй: transparency — operator должен "
        "знать, что система делает. Третий: SPOF analysis обязателен в FMEA/FTA. "
        "Четвёртый: software cannot solve hardware shortfalls.\n\n"
        "Patriot mini-callback. Friendly fire 2003 (RAF Tornado + USN F/A-18C, "
        "оба сбиты собственными Patriot) и украинский F-16 2024. Automation "
        "bias: когда система «лучше человека» по статистике, операторы "
        "перестают активно мониторить. Mitigation — системный подход, не single "
        "ML upgrade.")
    return s


def slide_30_act_criteria(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, "Act — когда не AI: 2 критерия")

    criteria = [
        ("user-check", "5",
         "Автономия не нужна; человек медленнее, но безопаснее",
         "MCAS — canonical контрпример. Auto-trim был «решением» проблемы, которой могло не быть",
         "СНАЧАЛА пересмотрите hardware. Software для compensation hardware shortfall — индикатор более глубокой проблемы"),
        ("hard-drive", "6",
         "COTS sensor дешевле и надёжнее, чем ML на одном sensor",
         "COTS = Commercial Off-The-Shelf. Не делайте ML на проблеме, решающейся hardware redundancy",
         "Второй AoA-сенсор на 737 MAX стоил бы порядки меньше всех trim-AI"),
    ]
    cc_w = 5.95; cc_h = 4.2; gap = 0.25
    for i, (icon, num, title, desc, alt) in enumerate(criteria):
        x = 0.6 + i * (cc_w + gap)
        ocean_box(s, x, 1.85, cc_w, cc_h)
        ocean_box(s, x + 0.3, 2.05, 0.7, 0.7,
                  fill=GOLD, stroke=GOLD, radius_pt=20)
        text_box(s, x + 0.3, 2.05, 0.7, 0.7, "#" + num,
                 size=16, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        icon_p = ASSETS / "icons" / f"{icon}-96.png"
        if icon_p.exists():
            add_image(s, icon_p, x + cc_w - 1.2, 2.05, w=0.9, h=0.9)
        text_box(s, x + 0.3, 3.0, cc_w - 0.6, 0.7, title,
                 size=16, bold=True, color=DEEP, line_spacing=1.2)
        text_box(s, x + 0.3, 3.85, cc_w - 0.6, 1.1, desc,
                 size=12, color=MID, italic=True, line_spacing=1.3)
        hr_line(s, x + 0.3, 5.05, cc_w - 0.6, color=LIGHT, weight=0.5)
        text_box(s, x + 0.3, 5.15, cc_w - 0.6, 0.95, "→ " + alt,
                 size=11, color=TEAL, italic=True, bold=True, line_spacing=1.3)

    ocean_box(s, 0.6, 6.35, 12.13, 0.65, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.43, 11.63, 0.5,
             "Act — звено, где hype далеко впереди реальности. Большинство strikes — operator-in-loop. Не путайте L3 с L5.",
             size=13, italic=True, color=DEEP, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_notes(s,
        "Из разобранных провалов извлекаем два критерия для звена Act — "
        "критерии 5 и 6.\n\n"
        "Критерий пятый — автономия не нужна. Канонический контрпример — "
        "MCAS: автоматическая корректировка триммирования была «решением» "
        "проблемы, которой могло не быть в первую очередь. Если ваша задача — "
        "auto-correct физическую проблему, сначала пересмотрите hardware. "
        "Software для compensation hardware shortfall — индикатор более "
        "глубокой проблемы.\n\n"
        "Критерий шестой — COTS sensor дешевле. COTS — Commercial Off-The-"
        "Shelf. Если проблема решается hardware redundancy, не делайте ML на "
        "проблеме, которая решается железом. Второй AoA-сенсор на 737 MAX "
        "стоил бы порядки меньше.\n\n"
        "Подытог Раздела 3. Act — звено, где hype далеко впереди реальности. "
        "Программы автономии растут быстро. Но большинство strikes остаётся "
        "operator-in-loop. Дальше — мета-уровень: где звено Act обрезано "
        "регулированием.")
    return s


# Placeholder builders for s31-43 (build out incrementally)
def slide_31(prs):
    return section_divider(prs, 4, "Граница и регулирование",
        "Где звено Act обрезано международным правом и государственной политикой",
        current_section=4,
        caption="15 минут · L1-L5 + UN GGE + ICRC + Maven shift + HITL/HOOL/HOTL · целиком strict-in")


def slide_32_l1_l5(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Лестница автономии L1-L5: что делает AI / что делает человек / ms-to-intervention",
        size=22)

    # Build 5-row table-style ladder
    levels = [
        ("L1", "Assistive",
         "AI: выдаёт detections · Человек: решает",
         "Palantir MSS analyst surface",
         "минуты-часы", LIGHT_TINT, LIGHT),
        ("L2", "Semi-auto perception",
         "AI: рекомендует action · Человек: авторизует",
         "Saker Scout target lock",
         "seconds", LIGHT_TINT, LIGHT),
        ("L3", "Supervised autonomy",
         "AI: executes в envelope · Человек: supervises",
         "Anduril Fury wingman (CCA)",
         "100-1000 ms", LIGHT_TINT, MID),
        ("L4", "Pre-authorised auto-engage",
         "AI: engages по ROE · Человек: может intervene",
         "Patriot auto, S-400 auto ROE",
         "<100 ms", TEAL_TINT, TEAL),
        ("L5", "Full LAWS",
         "AI: lethal без human · Человек: вне loop",
         "Currently debated, not deployed",
         "N/A — вне loop", GOLD_TINT, RED_WARN),
    ]
    row_h = 0.78; row_y = 1.85
    col_widths = [0.7, 2.2, 4.2, 3.0, 1.5]

    # Header
    headers = ["#", "Уровень", "AI / Человек", "Пример 2026", "ms-intervention"]
    cum_x = 0.6
    for i, h in enumerate(headers):
        filled_rect(s, cum_x, row_y, col_widths[i], 0.42, MID)
        text_box(s, cum_x + 0.08, row_y + 0.05, col_widths[i] - 0.16, 0.32, h,
                 size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        cum_x += col_widths[i]

    # Body rows — bottom-aligned conceptually (L5 on top of stack visually, L1 bottom)
    for r_i, (lvl, name, ai_human, ex, ms, bg, stroke) in enumerate(levels):
        y = row_y + 0.42 + r_i * row_h
        cum_x = 0.6
        cells = [lvl, name, ai_human, ex, ms]
        for c_i, cell in enumerate(cells):
            cell_color = bg
            txt_color = DEEP
            bold = False
            if c_i == 0:
                cell_color = stroke; txt_color = WHITE; bold = True
            elif c_i == 1:
                bold = True
            elif c_i == 4:
                cell_color = GOLD_TINT if r_i >= 3 else LIGHT_TINT
            filled_rect(s, cum_x, y, col_widths[c_i], row_h, cell_color,
                        stroke=LIGHT, stroke_pt=0.5)
            text_box(s, cum_x + 0.1, y + 0.08, col_widths[c_i] - 0.2, row_h - 0.16,
                     cell,
                     size=10 if c_i == 2 else 12,
                     bold=bold, color=txt_color, italic=(c_i == 3),
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
            cum_x += col_widths[c_i]

    # Boundary callouts on right side of L3 and L4-L5 boundary
    text_box(s, 12.3, row_y + 0.42 + 2*row_h, 1.0, 0.5,
             "← L3↔L4\nengineering\ndebate",
             size=8, italic=True, color=MID, bold=True, line_spacing=1.1)
    text_box(s, 12.3, row_y + 0.42 + 4*row_h, 1.0, 0.5,
             "← L4↔L5\ntreaty\ndebate",
             size=8, italic=True, color=GOLD, bold=True, line_spacing=1.1)

    # Bottom takeaway
    text_box(s, 0.6, 7.0, 12.13, 0.3,
             "Студент-инженер: сказать про конкретную систему НА КАКОМ УРОВНЕ. Не «автономная» — а «L3 с envelope шириной X»",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Чтобы говорить о границе «где можно, где нельзя», нужна общая шкала. "
        "В индустрии используется лестница L1-L5 — концептуальный аналог SAE-"
        "уровней для cars. К 2026 году L1-L5 сложилась как наиболее цитируемая "
        "операциональная шкала.\n\n"
        "Ключ — на каждом уровне явно сказано, что делает AI и что делает "
        "человек. Без этого «уровень» — пустое слово.\n\n"
        "L1 Assistive — AI выдаёт information, человек решает. Palantir MSS. "
        "ms-to-intervention — минуты-часы. L2 Semi-auto — AI рекомендует, "
        "человек авторизует. Saker Scout. Секунды. L3 Supervised — AI executes "
        "в pre-authorised envelope, человек supervises. Anduril Fury. 100-1000 "
        "ms. L4 Pre-authorised — AI engages по pre-set ROE, человек может "
        "intervene но не required. Patriot auto, S-400 auto. <100 ms. L5 Full "
        "LAWS — AI executes lethal без human authorisation. Currently debated, "
        "not deployed.\n\n"
        "Две границы. L3 ↔ L4 — место инженерного спора. Pre-authorisation "
        "envelope насколько узок? L4 ↔ L5 — место юридического спора в UN GGE. "
        "Даже Lavender формально требует human approval — 20-сек, — то есть "
        "формально L4-edge, не L5.\n\n"
        "Студент должен уметь сказать про систему: «эта система L3, с "
        "envelope шириной X». Это профессиональный язык в этой области.")
    return s


def slide_33_un_gge(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "UN GGE on LAWS — третья подряд резолюция 2025 (164/6/7 UN press)",
        size=22)

    # Left: vertical timeline
    ocean_box(s, 0.6, 1.95, 7.8, 4.7)
    text_box(s, 0.8, 2.1, 7.4, 0.45, "Хронология 2024-2025",
             size=15, bold=True, color=DEEP)

    timeline = [
        ("5 Nov 2024", "First Committee UNGA", "161 / 3 / 13", MID,
         "Против: Беларусь · КНДР · Россия"),
        ("2 Dec 2024", "Pleno UNGA · Резолюция 79/62", "166 / 3 / 15", MID,
         ""),
        ("Sep 2025", "UN GGE — 42 states joint statement", "rolling text", LIGHT,
         "объявлен достаточной основой для переговоров"),
        ("6 Nov 2025", "First Committee · третья подряд", "164 / 6 / 7", GOLD,
         "Против: Беларусь · Бурунди · КНДР · Израиль · Россия · США"),
    ]
    for i, item in enumerate(timeline):
        date, body, vote, col, note = item
        y = 2.6 + i * 1.0
        # Dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.95), Inches(y + 0.15),
                                 Inches(0.25), Inches(0.25))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background()
        text_box(s, 1.4, y, 1.6, 0.4, date,
                 size=12, bold=True, color=col)
        text_box(s, 3.05, y, 3.5, 0.35, body,
                 size=11, color=DEEP, italic=True)
        text_box(s, 6.6, y, 1.7, 0.4, vote,
                 size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)
        if note:
            text_box(s, 1.4, y + 0.45, 7.0, 0.4, note,
                     size=10, italic=True, color=MID, line_spacing=1.2)

    # Right: DoD Directive callout
    ocean_box(s, 8.5, 1.95, 4.3, 4.7, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 8.7, 2.1, 4.0, 0.45, "DoD Directive 3000.09",
             size=14, bold=True, color=DEEP)
    text_box(s, 8.7, 2.55, 4.0, 0.35,
             "«Autonomy in Weapon Systems»",
             size=11, italic=True, color=MID)
    text_box(s, 8.7, 2.95, 4.0, 0.35,
             "2012, updated 2023",
             size=11, italic=True, color=MID)

    hr_line(s, 8.7, 3.4, 3.8, color=TEAL, weight=0.5)

    text_box(s, 8.7, 3.55, 4.0, 0.5,
             "US policy формально требует HITL для kinetic engagement",
             size=11, color=DEEP, italic=True, line_spacing=1.25)

    text_box(s, 8.7, 4.25, 4.0, 0.4, "Системы по умолчанию:",
             size=11, bold=True, color=DEEP)
    text_runs(s, 8.7, 4.65, 4.0, 0.5, [
        {"text": "L1-L3", "size": 18, "bold": True, "color": GOLD},
        {"text": ", waiver для L4", "size": 12, "color": DEEP},
    ])

    text_box(s, 8.7, 5.5, 4.0, 1.1,
             "US не противник договора в принципе, но опасается жёстких ограничений. "
             "Сдвиг 2024→2025: «за» → «против»",
             size=10, italic=True, color=MID, line_spacing=1.3)

    # Goal callout
    ocean_box(s, 0.6, 6.85, 12.13, 0.4, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.85, 6.88, 11.63, 0.35,
             "Цель Генсека ООН: договор к 2026 [VFY-day-of]",
             size=13, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_notes(s,
        "UN GGE on LAWS — основной международный форум переговоров по LAWS. "
        "Создан в 2016. За последние три года произошёл сдвиг от «обсуждаем» к "
        "«работаем над текстом».\n\n"
        "Хронология. 5 ноября 2024 — Первый комитет UNGA: 161 за, 3 против, 13 "
        "воздержавшихся. Против — Беларусь, КНДР, Россия. 2 декабря 2024 — "
        "резолюция 79/62 на пленарном: 166 / 3 / 15. Сентябрь 2025 — UN GGE: "
        "42 государства joint statement; rolling text объявлен достаточной "
        "основой. 6 ноября 2025 — Первый комитет, третья подряд: 164 / 6 / 7 по "
        "UN press; 156 / 5 / 8 по Stop Killer Robots. Против — Беларусь, "
        "Бурунди, КНДР, Израиль, Россия, и теперь США. Это значимый сдвиг.\n\n"
        "DoD Directive 3000.09. Autonomy in Weapon Systems, 2012, updated 2023. "
        "US policy требует HITL для kinetic в большинстве сценариев. Системы по "
        "умолчанию в L1-L3, с явным waiver-процессом для L4. US не противник "
        "договора в принципе, но опасается жёстких ограничений.\n\n"
        "Цель Генсека ООН — договор к 2026. Если ваша система проектируется на "
        "горизонте 5-10 лет, она почти наверняка будет работать в рамках "
        "какого-то ограничения.")
    return s


def slide_34_icrc_skr(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "ICRC + Stop Killer Robots — non-state давление на переговоры", size=22)

    # Left: ICRC
    ocean_box(s, 0.6, 1.95, 6.0, 4.7)
    # Mini icon
    icon_p = ASSETS / "icons" / "shield-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 0.8, 2.1, w=0.5, h=0.5)
    text_box(s, 1.4, 2.15, 5.0, 0.4, "ICRC",
             size=18, bold=True, color=DEEP)
    text_box(s, 1.4, 2.55, 5.0, 0.3,
             "International Committee of the Red Cross",
             size=11, italic=True, color=MID)

    text_box(s, 0.8, 3.05, 5.6, 0.35, "Prohibit (полный запрет):",
             size=12, bold=True, color=DEEP)
    proh = ["непредсказуемые автономные системы",
            "AWS против людей"]
    for i, p in enumerate(proh):
        text_box(s, 0.95, 3.4 + i * 0.3, 5.4, 0.28, "• " + p,
                 size=11, color=MID, line_spacing=1.2)

    text_box(s, 0.8, 4.05, 5.6, 0.35,
             "Restrict (ограничения): остальные AWS",
             size=11, color=MID, italic=True)

    hr_line(s, 0.8, 4.5, 5.4, color=LIGHT, weight=0.5)

    # 2 quote boxes
    ocean_box(s, 0.8, 4.65, 5.55, 0.85, fill=LIGHT_TINT)
    text_box(s, 0.95, 4.75, 5.25, 0.7,
             "«Ceding life-and-death decisions to machines is dehumanizing»",
             size=11, italic=True, color=DEEP, line_spacing=1.25, bold=True)

    ocean_box(s, 0.8, 5.6, 5.55, 1.0, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.95, 5.7, 5.25, 0.85,
             "«It is not the weapon system that must comply with IHL, "
             "but the humans using it»",
             size=12, italic=True, color=DEEP, bold=True, line_spacing=1.25)

    # Right: Stop Killer Robots
    ocean_box(s, 6.8, 1.95, 5.95, 4.7)
    icon_p = ASSETS / "icons" / "users-48.png"
    if icon_p.exists():
        add_image(s, icon_p, 7.0, 2.1, w=0.5, h=0.5)
    text_box(s, 7.6, 2.15, 5.0, 0.4, "Stop Killer Robots",
             size=18, bold=True, color=DEEP)
    text_box(s, 7.6, 2.55, 5.0, 0.3,
             "Coalition · 270 NGOs · 70 countries",
             size=11, italic=True, color=MID)

    # Big 30 countries
    text_runs(s, 7.0, 3.05, 5.5, 0.7, [
        {"text": "30 стран", "size": 28, "bold": True, "color": GOLD},
        {"text": " явно поддерживают", "size": 13, "color": DEEP},
        {"newpara": True, "text": "полный запрет на fully autonomous weapons",
         "size": 11, "italic": True, "color": MID},
    ])

    # Country grid (compact)
    countries = "Алжир · Аргентина · Австрия · Боливия · Бразилия · Чили · Китай · Колумбия · Коста-Рика · Куба · Джибути · Эквадор · Египет · Эль-Сальвадор · Гана · Гватемала · Святой Престол · Ирак · Иордания · Мексика · Марокко · Намибия · Никарагуа · Пакистан · Панама · Перу · Палестина · Уганда · Венесуэла · Зимбабве"
    text_box(s, 7.0, 4.0, 5.55, 1.95, countries,
             size=9, color=MID, italic=True, line_spacing=1.4)

    text_box(s, 7.0, 6.05, 5.55, 0.55,
             "Leverage: public awareness · Slaughterbots video · правозащитная аналитика",
             size=10, italic=True, color=TEAL, line_spacing=1.25, bold=True)

    # Bridge bottom
    add_footer(s,
        "Источники: ICRC 2024 position paper · Vienna Conference 2024 · HRW 2020-2025")

    add_speaker_notes(s,
        "Помимо государств, важны два игрока non-state. ICRC — главный "
        "международный авторитет по применению международного гуманитарного "
        "права. Stop Killer Robots — коалиция 270 НКО из 70 стран.\n\n"
        "Позиция ICRC. Prohibit: непредсказуемые AWS; AWS против людей. "
        "Restrict: все остальные AWS.\n\n"
        "Этическое ядро: «Ceding life-and-death decisions to machine sensors "
        "and software is a dehumanizing process». Делегирование решений о жизни "
        "и смерти сенсорам и софту — дегуманизирующий процесс.\n\n"
        "Процедурное ядро: «It is not the weapon system that must comply with "
        "IHL, but the humans using it». Это не оружейная система должна "
        "соответствовать международному гуманитарному праву — это люди, "
        "использующие её. Я хочу, чтобы вы вынесли эту мысль в карман.\n\n"
        "Stop Killer Robots. К 2025 — 30 стран явно поддерживают полный запрет: "
        "страны Латинской Америки, Африки, ряд арабских, Святой Престол, Китай. "
        "Эти две организации — основное non-state pressure. Их goal — fast-track "
        "к binding treaty.")
    return s


def slide_35_maven_shift(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Maven walkout 2018 → vendor replacement → big-tech возврат 2024-2026",
        size=22)

    # 3 eras horizontal
    eras = [
        ("ERA 1: Maven walkout", "2018", LIGHT_TINT, LIGHT, "door-open",
         ["Март 2018: Google leak",
          "4 000+ подписей · ~12 резигнаций",
          "Июнь 2018: Google не продлевает"]),
        ("ERA 2: Vendor replacement", "2018-2024", LIGHT_TINT, MID, "building-2",
         ["Anduril $30,5 млрд",
          "Palantir $60B капитализация",
          "Scale, Helsing — растут"]),
        ("ERA 3: Big-tech возврат", "2024-2026", GOLD_TINT, GOLD, "refresh-ccw",
         ["Jan 2024: OpenAI · ban removed",
          "Nov 2024: Anthropic IL6 (pivot)",
          "Sep 2025: Google возвращается"]),
    ]
    ew = 4.0; gap = 0.1
    for i, (title, date, bg, col, icon, bullets) in enumerate(eras):
        x = 0.6 + i * (ew + gap)
        ocean_box(s, x, 1.95, ew, 4.4, fill=bg, stroke=col)
        # Icon
        icon_p = ASSETS / "icons" / f"{icon}-48.png"
        if icon_p.exists():
            add_image(s, icon_p, x + 0.2, 2.1, w=0.5, h=0.5)
        text_box(s, x + 0.8, 2.15, ew - 0.9, 0.4, title,
                 size=13, bold=True, color=DEEP, line_spacing=1.15)
        text_box(s, x + 0.8, 2.5, ew - 0.9, 0.3, date,
                 size=11, italic=True, color=col, bold=True)
        hr_line(s, x + 0.2, 3.0, ew - 0.4, color=col, weight=0.5)

        for j, b in enumerate(bullets):
            text_box(s, x + 0.25, 3.15 + j * 0.55, ew - 0.45, 0.5, "• " + b,
                     size=11, color=DEEP, line_spacing=1.25)

        # Arrow to next
        if i < 2:
            add_arrow(s, x + ew + gap - 0.05, 3.8, 0.1, 0.4, fill=GOLD)

    # Bottom callout
    ocean_box(s, 0.6, 6.45, 12.13, 0.65, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.5, 11.63, 0.55,
             "Personal ethics ≠ industry regulation. Только legal regulation "
             "(treaty) может блокировать adoption на indust-level",
             size=13, italic=True, color=DEEP, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s,
        "Источники: TechPolicy Press 2018 · Intercept 2024 · BusinessWire 2024 · CNBC 2025")

    add_speaker_notes(s,
        "С 2018 по 2026 произошёл сдвиг в позиции big-tech AI labs относительно "
        "оборонной работы. Короткая, но критически важная нарративная арка.\n\n"
        "Эра 1 — Maven walkout 2018. Google leak март 2018: помогает Пентагону. "
        "4 000+ подписей; ~12 резигнаций; staged walkouts. К июню 2018 Google "
        "не продлевает контракт.\n\n"
        "Эра 2 — vendor replacement 2018-2024. Anduril, Palantir, Scale, "
        "Helsing набирают вес. Маркетинговая позиция — «patriotic alternative». "
        "Anduril оценен в $30,5B; Palantir — $60B капитализация.\n\n"
        "Эра 3 — big-tech возврат 2024-2025. Январь 2024 — OpenAI тихо удалил "
        "ban на military use. Ноябрь 2024 — Anthropic-Palantir-AWS, Claude IL6. "
        "Сентябрь 2025 — Google возвращается через Google Cloud.\n\n"
        "Урок 1. Personal ethics не равно industry regulation. Внутренние "
        "этические policy одной фирмы не блокируют adoption military-AI на "
        "indust-level — только legal regulation может. Урок 2. «Не работать на "
        "DoD» теперь редкая роскошь. Студент 2026 года столкнётся с этим выбором "
        "и должен решать его осознанно.\n\n"
        "Между Maven walkout и Anthropic IL6 — всего 6 лет, и AI-индустрия "
        "прошла полный цикл от «отказа как принцип» до «военные контракты — "
        "критический revenue stream».")
    return s


def slide_36_hitl_hool_hotl(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "HITL / HOOL / HOTL — триада. Граница HOOL→HOTL = engineering decision",
        size=22)

    # 3 panel layout
    panels = [
        ("HITL", "Human-In-The-Loop", LIGHT_TINT, LIGHT,
         "Человек в КАЖДОЙ decision-point",
         "AI не действует без явной authorisation",
         "L1, L2",
         ["Palantir MSS analyst",
          "Saker Scout operator"]),
        ("HOOL", "Human-On-The-Loop", LIGHT_TINT, MID,
         "Человек SUPERVISES",
         "Может intervene, но не required",
         "L3, L4",
         ["Fury CCA wingman",
          "Patriot auto ROE"]),
        ("HOTL", "Human-Out-of-The-Loop", GOLD_TINT, RED_WARN,
         "Человек ВНЕ execution-loop",
         "Нет real-time intervention",
         "L5",
         ["Treaty-discussion",
          "Currently не deployed"]),
    ]
    pw = 4.0; gap = 0.1
    for i, item in enumerate(panels):
        abbr, full, bg, col, p1, p2, mapping, examples = item
        x = 0.6 + i * (pw + gap)
        ocean_box(s, x, 1.85, pw, 4.4, fill=bg, stroke=col)
        # Abbr badge
        text_box(s, x, 2.0, pw, 0.7, abbr,
                 size=36, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        text_box(s, x, 2.7, pw, 0.4, full,
                 size=13, italic=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        # Visual: stick figure + loop indicator
        # Simplified — use loop circle + position label
        loop = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x + (pw - 0.8) / 2), Inches(3.2),
                                  Inches(0.8), Inches(0.8))
        loop.fill.background()
        loop.line.color.rgb = col; loop.line.width = Pt(1.5)
        text_box(s, x, 3.25, pw, 0.7, "AI\nloop",
                 size=10, italic=True, color=col, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

        # Human stick figure relative to loop
        positions = {
            "HITL": (x + pw / 2 - 0.1, 3.4),  # внутри
            "HOOL": (x + pw / 2 - 0.1, 2.95),  # над loop
            "HOTL": (x + 0.4, 3.5),  # вне
        }
        hx, hy = positions[abbr]
        h_head = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(hx), Inches(hy),
                                    Inches(0.2), Inches(0.2))
        h_head.fill.solid(); h_head.fill.fore_color.rgb = DEEP
        h_head.line.fill.background()

        hr_line(s, x + 0.2, 4.15, pw - 0.4, color=col, weight=0.5)

        # Description
        text_box(s, x + 0.2, 4.3, pw - 0.4, 0.45, p1,
                 size=11, color=DEEP, line_spacing=1.2, bold=True)
        text_box(s, x + 0.2, 4.75, pw - 0.4, 0.45, p2,
                 size=10, color=MID, italic=True, line_spacing=1.2)

        # Mapping
        text_runs(s, x + 0.2, 5.25, pw - 0.4, 0.4, [
            {"text": "L1-L5: ", "size": 11, "color": DEEP},
            {"text": mapping, "size": 13, "bold": True, "color": col},
        ])

        # Examples
        for j, ex in enumerate(examples):
            text_box(s, x + 0.25, 5.65 + j * 0.27, pw - 0.45, 0.25, "• " + ex,
                     size=10, color=MID, italic=True, line_spacing=1.1)

    # Big engineering callout bottom
    ocean_box(s, 0.6, 6.4, 12.13, 0.75, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.85, 6.48, 11.63, 0.3,
             "Engineering takeaway: «сколько ms у оператора на intervention» = ФОРМАЛЬНАЯ категоризация",
             size=12, bold=True, color=DEEP)
    text_box(s, 0.85, 6.78, 11.63, 0.3,
             "10 сек = HOOL · 200 мс = формально HOOL, фактически HOTL · 5 мс = инженерно HOTL",
             size=11, italic=True, color=DEEP)
    add_speaker_notes(s,
        "Самая важная mental model этого раздела — триада уровней человеческого "
        "контроля. Та же лестница L1-L5, но взгляд с противоположной стороны: "
        "что делает не AI, а человек.\n\n"
        "HITL — Human-In-The-Loop. Человек в каждой decision-point. AI не "
        "действует без явной authorisation. Mapping: L1, L2. Примеры: Palantir "
        "MSS, Saker Scout operator. HOOL — Human-On-The-Loop. Человек "
        "supervises, может intervene, но не required в каждой точке. Mapping: "
        "L3, L4. Fury CCA wingman, Patriot auto. HOTL — Human-Out-of-The-Loop. "
        "Человек вне execution-loop, нет real-time intervention. Mapping: L5.\n\n"
        "Что важно для инженера. Граница HOOL → HOTL — это место, на которое "
        "заточены DoD Directive 3000.09 + UN GGE + ICRC. Граница формально "
        "определяется engineering decision: сколько ms у оператора на "
        "intervention. 10 секунд — HOOL. 200 мс — формально HOOL, фактически "
        "HOTL. 5 мс — инженерно HOTL.\n\n"
        "Engineering takeaway. «Сколько ms на intervention» — формальная "
        "категоризация системы с правовыми последствиями. Ответ должен быть "
        "зафиксирован в системных требованиях, не в маркетинге.\n\n"
        "Mitigation: calibrated uncertainty, abstention pathways, structured "
        "outputs, mandatory human gates для kinetic. Связь с провалами: "
        "Lavender — вырожденный HITL (20 сек → HOTL); MCAS — отсутствие "
        "meaningful override.")
    return s


def slide_37_russia_votes(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Россия против UN LAWS — инженер делает свой выбор внутри рамок",
        size=22)

    # Left: voting context
    ocean_box(s, 0.6, 1.95, 6.0, 4.7)
    text_box(s, 0.8, 2.1, 5.6, 0.4, "Voting context",
             size=16, bold=True, color=DEEP)

    votes = [
        ("Ноябрь 2024", "161 / 3 / 13", "Против: Беларусь · КНДР · Россия"),
        ("Ноябрь 2025", "164 / 6 / 7", "Против: + Бурунди · Израиль · США"),
    ]
    for i, (date, vote, against) in enumerate(votes):
        y = 2.7 + i * 1.55
        text_box(s, 0.95, y, 2.5, 0.4, date,
                 size=14, bold=True, color=MID if i == 0 else GOLD)
        text_box(s, 3.6, y, 2.7, 0.5, vote,
                 size=20, bold=True, color=DEEP, align=PP_ALIGN.RIGHT)
        text_box(s, 0.95, y + 0.5, 5.4, 0.85, against,
                 size=11, italic=True, color=MID, line_spacing=1.25)

    text_box(s, 0.8, 5.85, 5.6, 0.7,
             "Россия — в позиции «против» с 2018. Лагерь «против» в 2025 — больше 3 стран, политически разнообразен",
             size=11, color=DEEP, italic=True, line_spacing=1.3)

    # Right: 3 actions for engineer
    text_box(s, 6.9, 1.95, 5.9, 0.4, "Что делает инженер",
             size=16, bold=True, color=DEEP)

    actions = [
        ("book-open", "Знать ландшафт",
         "UN GGE, ICRC, голоса UNGA, DoD 3000.09. Как FAR для civil aviation"),
        ("scale", "Знать engineering определения",
         "HITL/HOOL/HOTL + L1-L5 применимы независимо от geopolitical alignment"),
        ("compass", "Делать осознанный выбор",
         "Civilian dual-use / оборона L1-L2 / оборона L3-L4 — каждый правомерный, но РАЗНЫЙ"),
    ]
    for i, (icon, title, desc) in enumerate(actions):
        y = 2.4 + i * 1.3
        ocean_box(s, 6.9, y, 5.9, 1.15)
        icon_p = ASSETS / "icons" / f"{icon}-48.png"
        if icon_p.exists():
            add_image(s, icon_p, 7.05, y + 0.15, w=0.55, h=0.55)
        text_box(s, 7.75, y + 0.15, 5.0, 0.4, f"{i+1}. {title}",
                 size=13, bold=True, color=DEEP)
        text_box(s, 7.75, y + 0.55, 5.0, 0.5, desc,
                 size=10, color=MID, italic=True, line_spacing=1.25)

    # Bottom takeaway
    ocean_box(s, 0.6, 6.55, 12.13, 0.55, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.85, 6.62, 11.63, 0.4,
             "В этой области инженер не остаётся нейтральным — он внутри рамок. Профессионализм — насколько осознанно их чувствует",
             size=13, italic=True, color=DEEP, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_notes(s,
        "Геополитический контекст для российского студента-инженера. Россия — "
        "одна из стран против резолюций UN LAWS. Ноябрь 2024: 161/3/13; против "
        "— Беларусь, КНДР, Россия. Ноябрь 2025: 164/6/7; против — Беларусь, "
        "Бурунди, КНДР, Израиль, Россия, США. Россия в позиции «против» с 2018. "
        "Лагерь «против» в 2025 — больше 3 стран, состав политически "
        "разнообразен.\n\n"
        "Мы не предписываем студенту политическую позицию — это его выбор. Но "
        "обязаны дать ландшафт.\n\n"
        "Первое — знать ландшафт. UN GGE, ICRC, голоса UNGA, DoD Directive — "
        "базовая профессиональная грамотность, как FAR для civil aviation "
        "engineer. Второе — знать, что критерии HITL/HOOL/HOTL и L1-L5 "
        "применимы независимо от geopolitical alignment. Engineering design "
        "определяется одинаково. Третье — делать осознанный выбор. Можно "
        "работать в гражданском dual-use, в обороне L1-L2, в обороне L3-L4. "
        "Каждый — правомерный, но разный.\n\n"
        "Это мысль, которую хотим донести: в этой области инженер не остаётся "
        "нейтральным — он внутри рамок, и профессионализм проявляется в том, "
        "насколько осознанно он эти рамки чувствует.\n\n"
        "Критерий 7, cross-cutting. Граница HOOL → HOTL — treaty-territory, "
        "не engineering. Когда проектируете систему рядом с этой границей, "
        "выходите за рамки чистой инженерии в зону международного права.")
    return s


def slide_38(prs):
    return section_divider(prs, 5, "Сборка: критерии, карьера, чтение, замыкание",
        "Семь критериев, карьерный угол, чтение, замыкание",
        current_section=5,
        caption="6 минут · 7 критериев · 5 профилей · 7 источников · замыкание к keystone")


def slide_39_seven_criteria(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Семь критериев «когда AI плохая идея» — рабочая матрица")

    criteria_data = [
        (1, "Sense", "Low-data domain или distribution shift", "Adversarial SAR ATR · новые цели", LIGHT),
        (2, "Sense", "High-stakes single-sensor без избыточности", "F-35 ALIS без HITL flight gate", LIGHT),
        (3, "Decide", "Long-tail edge cases с low ML confidence", "Mission planning под новые ROE", MID),
        (4, "Decide", "High-stakes life-and-death без HITL", "Lavender canonical anti-example", GOLD),
        (5, "Act", "Автономия не нужна, человек медленнее", "737 MAX MCAS — «решение не было нужно»", GOLD),
        (6, "Act", "COTS sensor дешевле + reliable", "AoA-redundancy на 737 MAX", TEAL),
        (7, "Cross-cutting", "HOOL → HOTL = treaty-territory", "LAWS · UN GGE", RED_WARN),
    ]
    col_widths = [0.5, 1.8, 5.5, 4.0]
    row_h = 0.55
    row_y = 1.85

    # Header
    headers = ["#", "Звено", "Критерий", "Иллюстрация"]
    cum_x = 0.6
    for i, h in enumerate(headers):
        filled_rect(s, cum_x, row_y, col_widths[i], 0.4, MID)
        text_box(s, cum_x + 0.1, row_y + 0.05, col_widths[i] - 0.2, 0.3, h,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cum_x += col_widths[i]

    # Body
    for r_i, (num, link, crit, illus, col) in enumerate(criteria_data):
        y = row_y + 0.4 + r_i * row_h
        bg = WHITE if r_i % 2 == 0 else SURFACE
        cum_x = 0.6
        is_highlight = num in [4, 5]
        if is_highlight:
            bg = GOLD_TINT
        cells = [str(num), link, crit, illus]
        for c_i, cell in enumerate(cells):
            filled_rect(s, cum_x, y, col_widths[c_i], row_h, bg,
                        stroke=LIGHT, stroke_pt=0.5)
            txt_color = col if c_i == 0 else (DEEP if c_i == 1 else MID)
            bold = (c_i == 0 or c_i == 1)
            italic = (c_i == 3)
            text_box(s, cum_x + 0.12, y + 0.08, col_widths[c_i] - 0.24,
                     row_h - 0.16, cell,
                     size=12, bold=bold, color=txt_color, italic=italic,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
            cum_x += col_widths[c_i]

    # Bottom callout
    ocean_box(s, 0.6, 6.6, 12.13, 0.5, fill=TEAL_TINT, stroke=TEAL)
    text_box(s, 0.85, 6.65, 11.63, 0.4,
             "Главное в матрице — она ИНСТРУМЕНТ. Один срабатывает → redesign; несколько → пересмотр подхода",
             size=13, italic=True, color=DEEP, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_notes(s,
        "Соберём в одну матрицу все шесть критериев плюс один cross-cutting — "
        "итого семь критериев «когда не AI» для аэрокосмической и оборонной "
        "области.\n\n"
        "В Sense — два. Первый: low-data domain или distribution shift. "
        "Иллюстрация — adversarial SAR ATR. Второй: high-stakes single-sensor "
        "без избыточности. Иллюстрация — F-35 ALIS без HITL.\n\n"
        "В Decide — два. Третий: long-tail edge cases с low ML confidence — "
        "нужна abstention. Четвёртый: high-stakes life-and-death без HITL. "
        "Канонический контрпример — Lavender.\n\n"
        "В Act — два. Пятый: автономия не нужна, человек медленнее. "
        "Канонический пример — 737 MAX MCAS как «решение проблемы, которой не "
        "было». Шестой: COTS sensor дешевле и надёжнее.\n\n"
        "Один cross-cutting, седьмой: граница HOOL → HOTL — это treaty-"
        "territory, не engineering. Это про LAWS и UN GGE.\n\n"
        "Главное в матрице — она инструмент. Прогоните предлагаемое AI-решение "
        "через эти семь критериев. Если хоть один срабатывает — это не «нельзя», "
        "это «нужен redesign». Если несколько — пересмотр фундаментального "
        "подхода. Эта матрица — не догма; она работает в 2026 с теми данными, "
        "которые есть. Рабочий инструмент.")
    return s


def slide_40_career(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s,
        "Карьерный угол — 5 инженерных профилей × 3 контура")

    # 5 profiles in a row
    profiles = [
        ("scan", "CV / DL",
         ["Спутниковая аналитика",
          "Defense ML",
          "Perception для drones"]),
        ("cpu", "ML / RL",
         ["Autonomous platforms",
          "Drone swarms",
          "Gen design"]),
        ("microchip", "Embedded / edge",
         ["On-orbit ML",
          "Drone autonomy",
          "On-platform inference"]),
        ("shield-check", "Safety eng",
         ["DO-178C / ARP4754A",
          "FMEA / FTA",
          "Redundancy design"]),
        ("scale", "Ethics / policy",
         ["UN GGE process",
          "ICRC engineering",
          "AI policy в orgs"]),
    ]
    pw = 2.4; gap = 0.1
    for i, (icon, title, items) in enumerate(profiles):
        x = 0.6 + i * (pw + gap)
        ocean_box(s, x, 1.85, pw, 2.55)
        icon_p = ASSETS / "icons" / f"{icon}-96.png"
        if icon_p.exists():
            add_image(s, icon_p, x + (pw - 0.9) / 2, 2.0, w=0.9, h=0.9)
        text_box(s, x, 3.0, pw, 0.45, title,
                 size=14, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        for j, it in enumerate(items):
            text_box(s, x + 0.15, 3.5 + j * 0.3, pw - 0.3, 0.28, "• " + it,
                     size=10, color=MID, italic=True, line_spacing=1.15)

    # 3 contours bottom
    text_box(s, 0.6, 4.65, 12.13, 0.35, "Три контура работы:",
             size=14, bold=True, color=DEEP)

    contours = [
        ("Российский dual-use", "Cognitive Pilot · VisionLabs · ТЕРРА ТЕХ · СКАНЭКС",
         LIGHT_TINT, LIGHT),
        ("Глобальный контур", "Boeing · Airbus · NASA · Anduril · Palantir · Helsing",
         LIGHT_TINT, MID),
        ("Гражданский без оборонной нагрузки",
         "Wisk Aero · Joby · Maxar · Planet · BlackSky", TEAL_TINT, TEAL),
    ]
    cw = 4.0; cgap = 0.1
    for i, (title, content, bg, col) in enumerate(contours):
        x = 0.6 + i * (cw + cgap)
        ocean_box(s, x, 5.1, cw, 1.5, fill=bg, stroke=col)
        text_box(s, x + 0.2, 5.2, cw - 0.4, 0.4, title,
                 size=12, bold=True, color=DEEP, line_spacing=1.15)
        text_box(s, x + 0.2, 5.65, cw - 0.4, 0.9, content,
                 size=11, color=MID, italic=True, line_spacing=1.25)

    # Bottom takeaway
    text_box(s, 0.6, 6.8, 12.13, 0.35,
             "Выбор есть — и он НЕ сводится к «либо военная индустрия, либо ничего»",
             size=13, italic=True, color=DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Без агитации — где в этой области реальные карьерные пути.\n\n"
        "Пять инженерных профилей. CV/DL specialist — спутниковая аналитика "
        "(Maxar/BlackSky), defense ML (Palantir, Helsing), perception для "
        "drones (Shield AI). ML/RL — autonomous platforms (CCA), drone swarms, "
        "gen design для aerospace components. Embedded/edge — on-orbit ML, "
        "drone autonomy, on-platform inference. Systems/safety — certification "
        "под DO-178C и ARP4754A, FMEA/FTA, redundancy design. Ethics/policy — "
        "UN GGE process, ICRC engineering, AI policy.\n\n"
        "Три контура. Российский dual-use: Cognitive Pilot, VisionLabs (под "
        "санкциями), ТЕРРА ТЕХ / СКАНЭКС / СПУТНИКС. Глобальный: Boeing, Airbus, "
        "NASA, ESA, Lockheed, Northrop, RTX; Anduril, Palantir, Helsing, Shield "
        "AI, Scale. Гражданский без оборонной нагрузки: Wisk Aero, Joby Aviation, "
        "Maxar, Planet, BlackSky.\n\n"
        "Главное — выбор есть, и он не сводится к «либо военная индустрия, "
        "либо ничего». Гражданская аэрокосмическая отрасль — полноценное "
        "профессиональное поле.")
    return s


def slide_41_reading_list(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, "Семь источников для тех, кто идёт дальше")

    readings = [
        ("book-open", "Scharre, P. (2018). Army of None",
         "W. W. Norton. Базовая книга по LAWS. На русском доступна. CNAS Director"),
        ("file-text", "CSIS / Bondar (Apr 2026)",
         "How Russia is Building a Sovereign Drone Ecosystem with AI-Driven Autonomy"),
        ("newspaper", "Abraham, Y. (2024). Lavender",
         "+972 / Local Call. Главная первичная публикация по Lavender"),
        ("shield", "ICRC Position Paper (2024)",
         "Autonomous Weapons Systems. Главный non-state документ. icrc.org"),
        ("plane", "DARPA ACE briefings 2023-2024",
         "Публичные материалы по X-62A VISTA. Что реально умеет AI-pilot"),
        ("scroll", "GAO-20-316 + GAO-22-105943",
         "F-35 ALIS / ODIN transition. US audit-документы"),
        ("users", "Stop Killer Robots briefs 2025",
         "Сводки по UN GGE и UNGA голосованиям. stopkillerrobots.org"),
    ]
    # 2 columns: 4 + 3
    n_left = 4; cw = 5.95; gap = 0.25
    for i, (icon, title, desc) in enumerate(readings):
        col = 0 if i < n_left else 1
        row = i if i < n_left else (i - n_left)
        x = 0.6 + col * (cw + gap)
        y = 1.85 + row * 1.15
        ocean_box(s, x, y, cw, 1.05)
        icon_p = ASSETS / "icons" / f"{icon}-48.png"
        if icon_p.exists():
            add_image(s, icon_p, x + 0.15, y + 0.15, w=0.45, h=0.45)
        text_box(s, x + 0.75, y + 0.12, cw - 0.85, 0.4, f"{i+1}. " + title,
                 size=12, bold=True, color=DEEP, line_spacing=1.15)
        text_box(s, x + 0.75, y + 0.55, cw - 0.85, 0.5, desc,
                 size=10, color=MID, italic=True, line_spacing=1.25)

    add_footer(s,
        "Все 104 источника главы — в разделе «Источники» в lec-09/chapter.md, сгруппированы по разделам")
    add_speaker_notes(s,
        "Для тех, кто хочет идти дальше — семь источников.\n\n"
        "Первый — книга Пола Скэра «Army of None» 2018, W. W. Norton. Базовая "
        "книга по LAWS, доступна на русском. Автор — директор CNAS.\n\n"
        "Второй — CSIS Bondar апрель 2026: «How Russia is Building a Sovereign "
        "Drone Ecosystem». Главный академический разбор. Открытый доступ.\n\n"
        "Третий — Abraham, +972 / Local Call 2024: «Lavender: The AI machine "
        "directing Israel's bombing spree in Gaza».\n\n"
        "Четвёртый — ICRC Position Paper on AWS 2024.\n\n"
        "Пятый — DARPA ACE Program briefings 2023-2024. Что реально умеет AI-"
        "pilot.\n\n"
        "Шестой — GAO-20-316 по F-35 ALIS и GAO-22-105943 по ODIN transition.\n\n"
        "Седьмой — Stop Killer Robots briefs 2025. Регулярные сводки.\n\n"
        "Все 104 источника главы — в разделе «Источники» в chapter.md, "
        "сгруппированы по разделам.")
    return s


def slide_42_closing_callback(prs):
    s = blank(prs); set_slide_bg(s, WHITE)
    add_assertion_title(s, "Цепь по-прежнему держит инженер")

    # Re-use OODA chain from s05 but compressed
    # Dual-use band
    filled_rect(s, 0.6, 1.6, 12.13, 0.25, SOFT_GREY)
    text_box(s, 0.6, 1.6, 12.13, 0.22,
             "Гражданское ↔ Военное · те же стеки, два контура",
             size=10, italic=True, color=DARK_GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    items = [
        ("Sense", "eye",
         "AI ускоряет → человек верифицирует ground truth",
         "(ALIS)"),
        ("Decide", "brain",
         "AI ассистирует → человек удерживает authority",
         "(Lavender — анти-пример)"),
        ("Act", "plane",
         "AI исполняет в envelope → человек supervises",
         "(CCA wingman, не replacement)"),
    ]
    card_w = 3.7; card_h = 2.7; gap = 0.55
    start_x = 0.6 + (12.13 - 3*card_w - 2*gap) / 2
    for i, (name, icon_name, payoff, ref) in enumerate(items):
        x = start_x + i * (card_w + gap)
        y = 2.1
        ocean_box(s, x, y, card_w, card_h)
        # Icon
        icon_p = ASSETS / "icons" / f"{icon_name}-96.png"
        if icon_p.exists():
            add_image(s, icon_p, x + (card_w - 0.7) / 2, y + 0.25, w=0.7, h=0.7)
        text_box(s, x, y + 1.05, card_w, 0.5, name,
                 size=22, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        text_box(s, x + 0.2, y + 1.6, card_w - 0.4, 0.7, payoff,
                 size=11, color=MID, italic=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)
        text_box(s, x + 0.2, y + 2.25, card_w - 0.4, 0.35, ref,
                 size=10, italic=True, color=LIGHT, bold=True,
                 align=PP_ALIGN.CENTER)
        # Arrows between cards
        if i < 2:
            add_arrow(s, x + card_w + 0.05, y + card_h / 2 - 0.2,
                      gap - 0.1, 0.4, fill=LIGHT)

    # Big gold callout
    ocean_box(s, 0.6, 5.05, 12.13, 1.35, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.85, 5.15, 11.63, 0.5,
             "Цепь по-прежнему держит инженер",
             size=32, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    text_box(s, 0.85, 5.7, 11.63, 0.6,
             "AI — инструмент в инженерных руках, не автономный субъект. "
             "Профессионализм — умение сказать «да» и «нет»",
             size=14, italic=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    # Forward
    text_box(s, 0.6, 6.7, 12.13, 0.35,
             "Следующие лекции: 10 — энергетика; 11 — транспорт и логистика. "
             "Цепь Sense → Decide → Act работает везде, где есть инженерное решение и физический мир",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Главная мысль, с которой мы хотим оставить студента: цепь по-прежнему "
        "держит инженер. AI вошёл в каждое звено, но он не заменил человека.\n\n"
        "Он ускорил Sense — но Sense без человеческой проверки не работает: "
        "ALIS, GPS-spoofing. Он ускорил Decide — но Decide без real HITL "
        "превращается в Lavender. Он расширил Act — но Act без supervised "
        "pilots overhead не выходит из demo-stage: X-62A, Lancet rollback.\n\n"
        "Те же CV/ML-pipelines, те же сенсорные стеки работают и в гражданской "
        "аэрокосмосе, и в обороне; инженер выбирает контур.\n\n"
        "Это не значит, что AI «маленький» или «обманчивый». Это значит, что "
        "в high-stakes отрасли AI — это инструмент в инженерных руках, а не "
        "автономный субъект. Профессионализм инженера в этой области — это "
        "умение сказать «да» там, где AI действительно даёт измеримое "
        "преимущество, и умение сказать «нет» там, где AI создаёт риск, "
        "который не закрывается одной моделью.\n\n"
        "Курс продолжается. В Лекции 10 — энергетика; в Лекции 11 — транспорт "
        "и логистика. И в каждой из них вы будете узнавать паттерны из этой "
        "главы — потому что цепь Sense → Decide → Act работает везде.\n\n"
        "Спасибо. Дальше — Q&A.")
    return s


def slide_43_qa(prs):
    s = blank(prs); set_slide_bg(s, WHITE)

    # Big Q&A
    text_box(s, 0.6, 1.5, 12.13, 3.0, "Q&A?",
             size=180, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # Icon
    icon_p = ASSETS / "icons" / "message-circle-question-96.png"
    if icon_p.exists():
        add_image(s, icon_p, (SLIDE_W_IN - 0.9) / 2, 4.65, w=0.9, h=0.9)

    # 3 backup prompts
    text_box(s, 0.6, 5.85, 12.13, 0.35,
             "Backup prompts (на случай, если зал молчит):",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER, bold=True)

    prompts = [
        "Почему США, Россия и Китай не подписывают LAWS-договор?",
        "+972 преувеличил про Lavender?",
        "Куда идти работать без оборонной нагрузки?",
    ]
    pw = 4.0; gap = 0.1
    for i, p in enumerate(prompts):
        x = 0.6 + i * (pw + gap)
        ocean_box(s, x, 6.3, pw, 0.7, fill=LIGHT_TINT)
        text_box(s, x + 0.15, 6.4, pw - 0.3, 0.5, p,
                 size=10, color=DEEP, italic=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    add_speaker_notes(s,
        "Открываем Q&A. У меня в запасе три типичных вопроса, которые могут "
        "возникнуть.\n\n"
        "Первый — про политический ландшафт. «Если ICRC и UN GGE так озабочены "
        "LAWS, почему США, Россия и Китай не подписывают договор?» Три причины. "
        "Военная — каждая инвестирует в автономию как стратегическое "
        "преимущество. Дипломатическая — в UN GGE нет single agreed determination "
        "LAWS. Техническая — даже если договор подписан, как проверить?\n\n"
        "Второй — критический. «Lavender обсуждается с осуждением. А что насчёт "
        "критики +972 — может, они преувеличили?» ЦАХАЛ официально опровергает. "
        "Сторона +972 опирается на 6 IDF intelligence officers. Independent "
        "reproducibility невозможна. Но ICRC и Lieber Institute разобрали кейс "
        "и нашли его серьёзным. Даже при максимально консервативной оценке — "
        "серьёзный кейс. Мы делаем engineering judgment на pattern «accuracy% "
        "as wrong metric» — pattern не зависит от того, кто делал ошибку.\n\n"
        "Третий — практический. «Куда идти работать без оборонной нагрузки?» "
        "Три направления: гражданская спутниковая аналитика (ТЕРРА ТЕХ, СКАНЭКС "
        "в РФ; Maxar, Planet в мире); гражданская авиация (Боинг, Airbus, "
        "eVTOL); гражданский dual-use транспорт (Cognitive Pilot, eVTOL).\n\n"
        "Открыт для ваших вопросов.")
    return s


# ========== Main builder for part2 ==========

def build_part2(prs):
    slide_09_constellation(prs)
    slide_10_edge_ai(prs)
    slide_11_russian_sat(prs)
    slide_12_predictive_maintenance(prs)
    slide_13_f35_alis(prs)
    slide_14_adversarial_gps(prs)
    slide_15_sense_criteria(prs)
    slide_16(prs)
    slide_17_decide_intro(prs)
    slide_18_palantir(prs)
    slide_19_scale_helsing(prs)
    slide_20_anthropic_russian(prs)
    slide_21_lavender(prs)
    slide_22_lancet_vincennes(prs)
    slide_23_decide_criteria(prs)
    slide_24(prs)
    slide_25_act_intro(prs)
    slide_26_fury(prs)
    slide_27_x62a_saker(prs)
    slide_28_geran_cognitive(prs)
    slide_29_mcas(prs)
    slide_30_act_criteria(prs)
    slide_31(prs)
    slide_32_l1_l5(prs)
    slide_33_un_gge(prs)
    slide_34_icrc_skr(prs)
    slide_35_maven_shift(prs)
    slide_36_hitl_hool_hotl(prs)
    slide_37_russia_votes(prs)
    slide_38(prs)
    slide_39_seven_criteria(prs)
    slide_40_career(prs)
    slide_41_reading_list(prs)
    slide_42_closing_callback(prs)
    slide_43_qa(prs)
