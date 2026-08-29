"""
apply_comments.py — applies editor comments #243-#248 to sem-01.pptx IN PLACE.

NOTE: build_sem01.py is stale (19-slide old numbering, see its header). It must NOT
be run. This script edits the CURRENT 20-slide rendered/sem-01.pptx directly, reusing
build_sem01's styling helpers (text_box, ocean_box, icon, chip, ...) for consistency.

Comments handled:
  #243 slide 2 — replace broken Material-Icon ligatures with real Lucide icon PNGs
  #244 slide 2 — add workain.ai project card (bottom row -> 3 cards)
  #245 (new)   — insert "Чат курса в MAX" slide after slide 4
  #246 slide 5 — remove YOLO hand-raise bar + round chips/captions; reframe as discussion
  #247 slide 6 — replace unreadable charts with big numbers (usage / tools / trust) + fresh data
  #248 slide 18 — base + 6 progressive-reveal quiz slides; delete slide 19 (memo)
"""
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import build_sem01 as B
from build_sem01 import (
    text_box, multipara_box, ocean_box, filled_rect, dashed_box, chip,
    icon, add_image, speaker_notes, disable_shadow,
    DEEP, MID, LIGHT, TEAL, SURFACE, WHITE, GOLD, SLATE, GOLD_TINT, SOFT_GREY,
)

HERE = Path(__file__).resolve().parent
PPTX = HERE / "sem-01.pptx"

prs = Presentation(str(PPTX))
S = prs.slides


def emu_in(v):
    return Emu(v).inches if v is not None else 0.0


# ============================================================
# Slide clone / reorder helpers
# ============================================================

def _copy_bg(src_slide, dest_slide):
    src_csld = src_slide._element.find(qn('p:cSld'))
    dst_csld = dest_slide._element.find(qn('p:cSld'))
    bg = src_csld.find(qn('p:bg'))
    if bg is not None:
        dst_csld.insert(0, copy.deepcopy(bg))


def _remap_media_rels(src_part, dest_part, spTree):
    for attr in (qn('r:embed'), qn('r:link')):
        for el in spTree.iter():
            rId = el.get(attr)
            if rId:
                rel = src_part.rels[rId]
                new_rId = dest_part.relate_to(rel.target_part, rel.reltype)
                el.set(attr, new_rId)


def duplicate_slide(index):
    """Deep-copy slide at `index`, append at end, return new slide."""
    src = S[index]
    dest = prs.slides.add_slide(src.slide_layout)
    # strip layout-inherited placeholder shapes
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    _copy_bg(src, dest)
    for shp in src.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    _remap_media_rels(src.part, dest.part, dest.shapes._spTree)
    return dest


def sldIdLst():
    return prs.slides._sldIdLst


def move_slide(new_pos, sldId_el):
    lst = sldIdLst()
    lst.remove(sldId_el)
    lst.insert(new_pos, sldId_el)


def delete_slide(index):
    lst = sldIdLst()
    sldId = list(lst)[index]
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    lst.remove(sldId)


# ============================================================
# Shape utilities
# ============================================================

def del_shape(shp):
    shp._element.getparent().remove(shp._element)


def set_only_run_text(shp, text):
    tf = shp.text_frame
    tf.paragraphs[0].runs[0].text = text


def add_fill(shp, color):
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


# ============================================================
# #243 + #244 — slide 2 (index 1): real icons + workain.ai card
# ============================================================

def fix_slide2():
    s = S[1]
    # Map: group top-Y band -> (icon name, hex, size relative to badge)
    # Contact badges (light E8F1F5) -> MID icons; project badges -> white/deep icons.
    # We locate each ligature textbox inside its group, clear it, and overlay a PNG.
    groups = [g for g in s.shapes if g.shape_type == 6]
    # Absolute badge geometry gathered earlier (inches):
    overlays = {
        'send':        ('send',        '065A82', 0.40, 6.17, 0.28),
        'mail':        ('mail',        '065A82', 0.40, 6.75, 0.28),
        'schema':      ('git-branch',  'FFFFFF', 4.66, 2.55, 0.36),
        'psychology':  ('brain-circuit','FFFFFF', 4.66, 4.22, 0.36),
        'shield':      ('shield-check','FFFFFF', 4.59, 5.68, 0.30),
        'history_edu': ('pencil-line', '21295C', 9.02, 5.68, 0.30),
    }
    # 1) clear ligature text inside groups
    for g in groups:
        for child in g.shapes:
            if child.has_text_frame:
                t = child.text_frame.text.strip()
                if t in overlays:
                    for para in child.text_frame.paragraphs:
                        for r in para.runs:
                            r.text = ""
    # 2) overlay real icon PNGs at absolute badge positions
    for key, (name, hexc, x, y, sz) in overlays.items():
        px = 96 if hexc == 'FFFFFF' and name in ('qr-code',) else 64
        icon(s, name, hexc, px, x, y, sz)

    # ---- #244: rebuild bottom row into THREE small cards (w4check, workain, blog)
    # remove the two existing small-card groups (w4check ~x4.27, blog ~x8.70, top 5.36)
    for g in list(groups):
        top = emu_in(g.top)
        if 5.2 < top < 5.5 and emu_in(g.width) < 5.0:
            del_shape(g)
    row_y, row_h = 5.36, 1.72
    gap = 0.26
    cw = (12.81 - 4.27 - 2 * gap) / 3     # ≈ 2.67
    x0 = 4.27
    cards = [
        ('shield-check', '065A82', 'w4check.io',
         'Проверка надёжности контрагентов в криптовалютах.', False),
        ('briefcase', '065A82', 'workain.ai',
         'AI-автоматизация рабочих процессов и задач.', False),
        ('pencil-line', 'F0AB00', 'Блог · tellian.io',
         'Авторские материалы о технологиях, ИИ и разработке.', True),
    ]
    for i, (icn, badge_hex, title, desc, gold) in enumerate(cards):
        cx = x0 + i * (cw + gap)
        box = ocean_box(s, cx, row_y, cw, row_h,
                        fill=(GOLD_TINT if gold else SURFACE),
                        stroke=(GOLD if gold else LIGHT))
        # icon badge
        bsz = 0.52
        badge = filled_rect(s, cx + 0.20, row_y + 0.20, bsz, bsz,
                            RGB(badge_hex), radius=True, radius_adj=0.28)
        ic = 'pencil-line' if icn == 'pencil-line' else icn
        ihex = '21295C' if badge_hex == 'F0AB00' else 'FFFFFF'
        icon(s, ic, ihex, 64, cx + 0.20 + bsz*0.18, row_y + 0.20 + bsz*0.18, bsz*0.62)
        text_box(s, cx + 0.20, row_y + 0.80, cw - 0.40, 0.30, title,
                 size=13.5, bold=True, color=DEEP)
        text_box(s, cx + 0.20, row_y + 1.12, cw - 0.40, row_h - 1.18, desc,
                 size=10.5, color=SLATE, line_spacing=1.1)


from pptx.dml.color import RGBColor
def RGB(hexc):
    return RGBColor(int(hexc[0:2], 16), int(hexc[2:4], 16), int(hexc[4:6], 16))


def set_lines(shp, lines):
    """Rewrite a text frame's first runs per paragraph, preserving formatting."""
    paras = shp.text_frame.paragraphs
    for i, line in enumerate(lines):
        if i < len(paras) and paras[i].runs:
            paras[i].runs[0].text = line
            for r in paras[i].runs[1:]:
                r.text = ''


# ============================================================
# Failures slide (orig index 6 = s07, becomes slide 8): instructor stories in notes
# ============================================================

FAILURE_STORIES = (
    "\n\n— Личные истории преподавателя (рассказать в этом блоке, по одной на категорию) —\n\n"
    "1. Галлюцинация факта. В прошлом году, когда я писал статью для конференции, модель "
    "выдала несколько ссылок на научные работы, которых не существует, — а часть ссылок вела "
    "на реальные статьи, но совсем по другой теме, чем она заявляла. Выглядело уверенно и "
    "правдоподобно, поэтому каждую ссылку пришлось проверять руками.\n\n"
    "2. Устаревшие данные. В разработке это особенно заметно: модель тянет устаревшие версии "
    "библиотек и, хуже того, устаревшие подходы. Например, просишь компонент на React — а она "
    "пишет его на классах с componentWillMount вместо хуков; или предлагает moment.js, который "
    "давно не развивается, вместо date-fns или встроенного Temporal. Код формально рабочий, но "
    "подход — из позавчерашнего дня.\n\n"
    "3. Задача решена только внешне. Так происходит почти с каждым продуктом, который я делаю: "
    "первый прототип от AI выглядит очень похоже на нужное — экран, кнопки, ответы, — но "
    "копнёшь внутрь, и часто оно либо не работает, либо имитирует работу (данные захардкожены, "
    "обработка ошибок нарисована, но не подключена). Снаружи готовый продукт, внутри декорация.\n\n"
    "4. Потеря контекста. В длинном диалоге по рефакторингу через пару десятков сообщений "
    "модель забывает то, о чём мы договорились в начале: скажем, я явно сказал «не тянем внешних "
    "зависимостей» и мы переименовали модуль — а через двадцать реплик она снова предлагает "
    "стороннюю библиотеку и старое имя, будто разговора не было. Контекст уехал за пределы окна, "
    "и условия приходится напоминать заново.\n\n"
    "5. Неуместный тон. Пример прямо перед вами — эти лекции. Я готовлю их с активным участием "
    "AI, и до сих пор на слайдах и в заметках можно поймать странные обороты и тон, которые я не "
    "дочистил. Это честный пример того, что тон модели надо править вручную.\n\n"
    "6. Ответ не под вашу ситуацию. У меня постоянно всплывает одна и та же история — "
    "параноидальная «информационная безопасность» модели на дев- и тест-стендах: она "
    "сопротивляется сохранить или показать логины и пароли тестовых пользователей, потому что "
    "«это небезопасно», хотя речь про заведомо тестовые учётки на локальном стенде, где это ровно "
    "то, что нужно. Совет правильный вообще — но не под мою конкретную ситуацию."
)


def fix_failures():
    s = S[6]
    # reframe category 6 label per owner instruction (comment follow-up)
    for shp in s.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip().startswith('Слишком общий совет'):
            set_lines(shp, ['Ответ не под вашу', 'конкретную ситуацию'])
    # append instructor stories to speaker notes
    cur = s.notes_slide.notes_text_frame.text if s.has_notes_slide else ''
    speaker_notes(s, cur.rstrip() + FAILURE_STORIES)


# ============================================================
# #246 — slide 5 (index 4): drop YOLO bar + rounds; reframe
# ============================================================

def fix_slide5():
    s = S[4]
    # title = topmost textbox (identify by position; keep it, never delete)
    title = next(sh for sh in s.shapes if sh.has_text_frame and emu_in(sh.top) < 0.7)
    to_del = []
    for shp in s.shapes:
        if shp._element is title._element:
            continue
        y = emu_in(shp.top)
        txt = shp.text_frame.text.strip() if shp.has_text_frame else ''
        fill = ''
        try:
            if shp.shape_type == 1 and shp.fill.type is not None:
                fill = str(shp.fill.fore_color.rgb)
        except Exception:
            pass
        # YOLO explainer bar (bg + 2 pics + 2 text lines), all above y=2.6
        if y < 2.6:
            to_del.append(shp); continue
        # round chips (teal, contain 'раунд')
        if fill == '028090' and 'раунд' in txt:
            to_del.append(shp); continue
        # round captions
        if txt.startswith('раунды поднятия руки'):
            to_del.append(shp); continue
    for shp in to_del:
        del_shape(shp)
    # retitle
    set_only_run_text(title, 'Разминка: как вы используете AI?')
    # drop lingering round-mechanic wording in an option line
    for shp in s.shapes:
        if shp.has_text_frame and 'раунд на каждое значение' in shp.text_frame.text:
            set_only_run_text(shp, 'шкала от 1 до 5')
    # subtitle in freed space
    text_box(s, 0.55, 1.45, 12.23, 0.55,
             'Отвечаем вслух и обсуждаем — без камеры и подсчёта. Ваши ответы мы сравним '
             'с внешней статистикой в следующем блоке.',
             size=15, color=SLATE, line_spacing=1.15)
    speaker_notes(s, SLIDE5_NOTES)


SLIDE5_NOTES = (
    "Небольшая разминка перед статистикой — знакомимся с тем, как группа работает с AI. "
    "Формат простой: я задаю вопрос, отвечаете вслух, при желании поднимаете руку — "
    "это не голосование с камерой и подсчётом, а живой разговор, чтобы увидеть картину группы.\n\n"
    "Четыре вопроса. Первый — какой AI-инструмент вы используете чаще всего: ChatGPT, "
    "GigaChat или YandexGPT, GitHub Copilot / Claude / DeepSeek, или не пользуетесь вовсе. "
    "Второй — как часто: каждый день, несколько раз в неделю, реже, никогда. Третий — "
    "насколько доверяете точности AI по шкале от одного до пяти. Четвёртый — был ли случай, "
    "когда AI вас явно подвёл.\n\n"
    "Эти ответы понадобятся в следующем блоке — будем сравнивать вашу группу с внешними "
    "исследованиями: выше вы или ниже по доверию, чаще или реже пользуетесь."
)


# ============================================================
# #247 — slide 6 (index 5): big numbers instead of charts
# ============================================================

def fix_slide6():
    s = S[5]
    for shp in list(s.shapes):
        del_shape(shp)
    # Title
    text_box(s, 0.55, 0.45, 12.30, 0.95, 'Используют чаще, доверяют меньше — и это не парадокс',
             size=28, bold=True, color=DEEP)

    def big_card(x, eyebrow, body_fn):
        cw = 3.87
        ocean_box(s, x, 1.62, cw, 3.55, fill=SURFACE, stroke=LIGHT)
        text_box(s, x + 0.28, 1.84, cw - 0.5, 0.30, eyebrow, size=12.5, bold=True, color=TEAL)
        body_fn(x, cw)

    # Card A — usage
    def usage(x, cw):
        text_box(s, x + 0.24, 2.20, cw - 0.4, 1.0, '84%', size=66, bold=True, color=MID)
        text_box(s, x + 0.28, 3.45, cw - 0.5, 0.7,
                 'используют или планируют использовать AI-инструменты', size=13, color=DEEP,
                 line_spacing=1.1)
        text_box(s, x + 0.28, 4.55, cw - 0.5, 0.5, '2024 → 2025:  76% → 84%  (+8 п.п.)',
                 size=12, bold=True, color=SLATE)
    big_card(0.55, 'ИСПОЛЬЗОВАНИЕ', usage)

    # Card B — tools breakdown (mini bars)
    def tools(x, cw):
        tb = [('ChatGPT', 82), ('GitHub Copilot', 68), ('Gemini', 47), ('Claude', 41)]
        ty = 2.28
        barx = x + 0.28
        barw = cw - 0.56
        for name, pct in tb:
            text_box(s, barx, ty, barw - 0.9, 0.28, name, size=12.5, bold=True, color=DEEP)
            text_box(s, barx + barw - 0.9, ty, 0.9, 0.28, f'{pct}%', size=12.5, bold=True,
                     color=MID, align=PP_ALIGN.RIGHT)
            # track + fill
            filled_rect(s, barx, ty + 0.30, barw, 0.13, SOFT_GREY, radius=True, radius_adj=0.5)
            filled_rect(s, barx, ty + 0.30, barw * pct / 100.0, 0.13, TEAL, radius=True, radius_adj=0.5)
            ty += 0.63
        text_box(s, barx, ty + 0.02, barw, 0.4, '% среди тех, кто пользуется AI',
                 size=10, italic=True, color=SLATE)
    big_card(4.72, 'ЧТО ИСПОЛЬЗУЮТ ЧАЩЕ ВСЕГО', tools)

    # Card C — trust
    def trust(x, cw):
        text_box(s, x + 0.24, 2.20, cw - 0.4, 1.0, '29%', size=66, bold=True, color=GOLD)
        text_box(s, x + 0.28, 3.45, cw - 0.5, 0.7,
                 'доверяют точности того, что выдаёт AI', size=13, color=DEEP, line_spacing=1.1)
        multipara_box(s, x + 0.28, 4.35, cw - 0.5, 0.8, [
            {'text': 'было 40% в 2024  (−11 п.п.)', 'size': 12, 'bold': True, 'color': SLATE},
            {'text': '46% явно не доверяют · лишь 3% доверяют полностью', 'size': 11,
             'color': SLATE, 'line_spacing': 1.1},
        ])
    big_card(8.89, 'ДОВЕРИЕ', trust)

    # Bottom strip — Russia (VCIOM) + sources
    ocean_box(s, 0.55, 5.35, 12.23, 1.15, fill=GOLD_TINT, stroke=GOLD)
    text_box(s, 0.78, 5.48, 3.4, 0.9, 'Россия · ВЦИОМ, 2026', size=12.5, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE)
    multipara_box(s, 4.2, 5.46, 8.3, 0.95, [
        {'text': '78% пользовались нейросетями за год (было 73%) · 58% — не реже раза в неделю '
                 '(было 51%)', 'size': 12.5, 'bold': True, 'color': DEEP, 'line_spacing': 1.12},
        {'text': '64% верят в пользу AI, из них 67% — только в отдельных сферах · '
                 'N=3209, 18+, 25–27.06.2026, ≤±1,7%', 'size': 10.5, 'color': SLATE,
         'line_spacing': 1.1},
    ], anchor=MSO_ANCHOR.MIDDLE)
    # sources line
    text_box(s, 0.55, 6.62, 12.23, 0.3,
             'Источники: Stack Overflow Developer Survey 2025 (N=49 009, май–июнь 2025; '
             'самоотобранная онлайн-выборка) · ВЦИОМ «Нейросети в нашей жизни»',
             size=9, italic=True, color=SLATE)
    speaker_notes(s, SLIDE6_NOTES)


SLIDE6_NOTES = (
    "Два независимых среза. Первый — Stack Overflow Developer Survey 2025, почти 49 тысяч "
    "разработчиков. Восемьдесят четыре процента используют или планируют использовать AI — "
    "рост с семидесяти шести процентов годом раньше. Чаще всего берут ChatGPT (82%) и "
    "GitHub Copilot (68%), заметно меньше — Gemini (47%) и Claude (41%). "
    "При этом доверяют точности вывода лишь двадцать девять процентов — падение с сорока; "
    "сорок шесть процентов прямо не доверяют, и только три процента доверяют полностью.\n\n"
    "Как так — пользуются чаще, а доверяют меньше? Это не парадокс, а признак зрелого "
    "использования: инструмент берут как черновик, который обязательно проверяют. "
    "Важна методология — это самоотобранный онлайн-опрос, аудитория смещена к англоязычным "
    "разработчикам.\n\n"
    "Второй срез — ВЦИОМ, репрезентативная выборка россиян. За год нейросетями пользовались "
    "семьдесят восемь процентов (было семьдесят три), еженедельно — пятьдесят восемь (было "
    "пятьдесят один). Здесь речь про население в целом, а не про инженеров.\n\n"
    "Сравните с вашими ответами из разминки: где ваша группа — выше или ниже по доверию, "
    "чаще или реже пользуется?"
)


# ============================================================
# #248 — slide 18 (index 17): base + 6 progressive reveals; delete s19
# ============================================================

# correct answers per question (1..6): True = 'верно', False = 'неверно'
CORRECT = {1: False, 2: False, 3: True, 4: False, 5: False, 6: False}

# Один развёрнутый комментарий на слайд — про ответ, открытый ИМЕННО на этом слайде
# (без повторения разбора предыдущих вопросов).
Q_NOTES = {
    1: ("Вопрос 1: «AI считает буквы в слове напрямую, как человек» — НЕВЕРНО.\n\n"
        "Модель не видит отдельные буквы. Перед обработкой текст разбивается на токены — куски "
        "из нескольких символов (иногда это целое слово, иногда часть слова или пара букв). Дальше "
        "модель работает уже с этими токенами, а не с буквами. Поэтому классический пример-провал — "
        "«сколько букв р в слове „характеристика“»: модель может ответить неверно, потому что она "
        "буквально не «считает буквы», а прикидывает ответ по статистике. Тот же корень у ошибок в "
        "подсчёте символов, переворачивании слов, задачах на анаграммы.\n\n"
        "Урок: там, где важен точный посимвольный результат (подсчёт, маски, форматы), не "
        "доверяйте модели на слово — проверьте или поручите это обычному коду, а не LLM."),
    2: ("Вопрос 2: «AI всегда знает, что произошло в мире сегодня» — НЕВЕРНО.\n\n"
        "У модели есть дата среза обучения: она знает мир примерно до этого момента и по умолчанию "
        "не в курсе того, что случилось позже — свежий релиз, новость, курс валют, вчерашний матч. "
        "Сам по себе чат не ходит в интернет. Чтобы ответ учитывал актуальные данные, нужен явный "
        "инструмент: веб-поиск, подключение к базе, загрузка свежего документа (это называют RAG). "
        "Без него уверенный ответ про «сегодня» — это в лучшем случае экстраполяция, в худшем — "
        "выдумка.\n\n"
        "Урок: для всего, что зависит от свежести (цены, версии, события), либо включайте поиск, "
        "либо проверяйте дату — «уверенно» не значит «актуально»."),
    3: ("Вопрос 3: «AI иногда уверенно говорит неправду, не намереваясь обмануть» — ВЕРНО.\n\n"
        "Это галлюцинация — мы уже называли её в блоке про провалы. Модель не «врёт» в человеческом "
        "смысле: у неё нет намерения обмануть, она просто достраивает наиболее правдоподобное "
        "продолжение и с той же интонацией уверенности выдаёт как факт, так и вымысел. Отсюда "
        "несуществующие ссылки, выдуманные цитаты, правдоподобные, но ложные детали. Ключевая "
        "ловушка именно в тоне: уверенность формулировки никак не связана с её правильностью.\n\n"
        "Урок: тон — не индикатор истины. Любой проверяемый факт (ссылка, цифра, цитата, API) "
        "сверяйте с первоисточником, особенно когда ответ звучит гладко и уверенно."),
    4: ("Вопрос 4: «AI-чат по умолчанию видит вашу личную почту и файлы» — НЕВЕРНО.\n\n"
        "Обычный чат видит только то, что вы сами прислали в диалоге. Он не читает вашу почту, "
        "диск или файлы на компьютере, пока вы явно не подключите такую интеграцию или не загрузите "
        "документ. Доступ появляется только там, где вы его осознанно дали, — и вот об этом уже "
        "стоит думать: что именно вы вставляете в промпт и какие интеграции включаете, особенно с "
        "рабочими данными.\n\n"
        "Урок: паниковать про «он всё видит» не нужно, но и вставлять секреты/персональные данные "
        "в чат по привычке — тоже; доступ определяете вы, значит вы и отвечаете за то, что даёте."),
    5: ("Вопрос 5: «AI всегда выдаёт ответ в соответствии с данными, которые в него заложили» — "
        "НЕВЕРНО.\n\n"
        "Обучающие данные формируют «характер» модели и её систематические искажения, но ответ не "
        "выводится из них жёстко и однозначно. Во-первых, модель не ищет готовый кусок текста в "
        "выборке, а генерирует продолжение по выученным статистическим закономерностям — поэтому "
        "возможны и обобщения на новые случаи, и галлюцинации. Во-вторых, из-за случайности "
        "сэмплирования (temperature) один и тот же вопрос при разных запусках даёт разные ответы. "
        "В-третьих, поведение сильно правится уже после предобучения — дообучением и настройкой под "
        "диалог (RLHF), так что финальный ответ отражает не только «сырые» тренировочные данные.\n\n"
        "Урок: не думайте о модели как о поисковике по базе. Это вероятностный генератор — отсюда и "
        "гибкость, и невоспроизводимость, и необходимость проверки."),
    6: ("Вопрос 6: «Все AI-модели дают одинаковый ответ на один и тот же вопрос» — НЕВЕРНО.\n\n"
        "Разные модели обучены по-разному и отвечают по-разному; более того, даже одна и та же "
        "модель при повторных запусках может дать разные по содержанию ответы (см. предыдущий пункт "
        "про случайность генерации). Это не баг, а свойство технологии. Практический плюс: если по "
        "важному вопросу задать его двум-трём разным моделям и ответы сходятся — это слабый, но "
        "полезный сигнал; если расходятся — повод копать глубже и идти к первоисточнику.\n\n"
        "Урок: для значимых решений сверяйтесь с несколькими источниками и, по возможности, с "
        "несколькими моделями — единственный ответ единственной модели не является истиной в "
        "последней инстанции."),
}

BASE_NOTE = ("Блиц-квиз: я зачитываю утверждение, вы голосуете поднятием руки — верно или "
             "неверно — прежде чем я дам ответ. На этом слайде ответов ещё нет; дальше правильные "
             "ответы открываются по одному, с коротким объяснением и уроком.")


def chip_qindex(top_in):
    return round((top_in - 1.82) / 0.94) + 1


def highlight_reveal(slide, upto_k):
    """On `slide`, mark correct chips for questions 1..upto_k."""
    # gather chip textboxes
    chips = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            t = shp.text_frame.text.strip()
            if t in ('верно', 'неверно'):
                chips.append(shp)
    for shp in chips:
        q = chip_qindex(emu_in(shp.top))
        if q < 1 or q > 6 or q > upto_k:
            continue
        label = shp.text_frame.text.strip()
        is_correct = (label == 'верно') == CORRECT[q]
        run = shp.text_frame.paragraphs[0].runs[0]
        if is_correct:
            add_fill(shp, GOLD)
            run.text = '✓ ' + label
            run.font.bold = True
            run.font.color.rgb = WHITE
        else:
            run.font.color.rgb = SOFT_GREY
            run.font.bold = False
    # update small print
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip().startswith('голосуем'):
            set_only_run_text(shp, f'правильные ответы открываются по одному — вопрос {upto_k} из 6')


def build_reveal_slides():
    base_idx = 17           # slide 18 (quiz base)
    # update base slide small print + notes
    base = S[base_idx]
    speaker_notes(base, BASE_NOTE)
    # create 6 reveals (append at end), remember their sldIds in order
    new_sldids = []
    for k in range(1, 7):
        dup = duplicate_slide(base_idx)
        highlight_reveal(dup, k)
        # one expanded comment per slide — only the answer revealed here, no repeats
        speaker_notes(dup, Q_NOTES[k])
        # its sldId is the last one in the list
        new_sldids.append(list(sldIdLst())[-1])
    return new_sldids


# ============================================================
# #245 — new "Чат курса в MAX" slide (built fresh, appended, then moved)
# ============================================================

def build_max_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    B.set_slide_bg(s, WHITE)
    text_box(s, 0.55, 0.45, 12.30, 0.8, 'Чат курса в MAX — всё важное в одном месте',
             size=28, bold=True, color=DEEP)
    text_box(s, 0.55, 1.30, 12.30, 0.5, 'Чат группы: «Отраслевое использование ИИ 2026»',
             size=17, bold=True, color=TEAL)

    # Left: real QR to the MAX course chat + the link underneath
    qx, qy, qw, qh = 0.9, 1.95, 4.0, 4.35
    ocean_box(s, qx, qy, qw, qh, fill=WHITE, stroke=LIGHT, stroke_pt=1.8)
    qsz = 2.55
    add_image(s, B.SHOTS / 's04b-max-qr.png', qx + (qw - qsz) / 2, qy + 0.28, w=qsz, h=qsz)
    text_box(s, qx + 0.2, qy + 2.98, qw - 0.4, 0.4, 'Отсканируйте камерой телефона',
             size=14, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    text_box(s, qx + 0.15, qy + 3.44, qw - 0.3, 0.8,
             'max.ru/join/sHoHlhI4jvW_swcq5ZLowvz6G94m0IMhAwngxuxsEpI',
             size=10, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.12)

    # Right: what's in the chat
    rx = 5.4
    rows = [
        ('bell', 'Объявления', 'Расписание, изменения, важные новости курса — в одном месте.'),
        ('book-open', 'Материалы лекций и семинаров',
         'Слайды и материалы появляются в чате по мере проведения занятий.'),
        ('message-square-quote', 'Вопросы между занятиями',
         'Можно задать вопрос по курсу и получить ответ, не дожидаясь пары.'),
    ]
    ry = 1.95
    rh = 1.35
    gap = 0.15
    for icn, title, desc in rows:
        ocean_box(s, rx, ry, 7.4, rh, fill=SURFACE, stroke=LIGHT)
        bsz = 0.62
        filled_rect(s, rx + 0.24, ry + 0.30, bsz, bsz, TEAL, radius=True, radius_adj=0.26)
        icon(s, icn, 'FFFFFF', 64, rx + 0.24 + bsz*0.18, ry + 0.30 + bsz*0.18, bsz*0.62)
        text_box(s, rx + 1.15, ry + 0.24, 7.4 - 1.4, 0.4, title, size=16, bold=True, color=DEEP)
        text_box(s, rx + 1.15, ry + 0.68, 7.4 - 1.4, 0.55, desc, size=12.5, color=SLATE,
                 line_spacing=1.12)
        ry += rh + gap

    speaker_notes(s, MAX_NOTE)
    return list(sldIdLst())[-1]


MAX_NOTE = (
    "Всё общение по курсу — в чате в мессенджере MAX. Отсканируйте QR-код камерой телефона "
    "или перейдите по ссылке max.ru/join/sHoHlhI4jvW_swcq5ZLowvz6G94m0IMhAwngxuxsEpI "
    "и вступите в чат до первой лекции.\n\n"
    "Что там будет. Во-первых, объявления: расписание, любые изменения, важные новости — "
    "всё публикуется в чате, отдельно письма ждать не нужно. Во-вторых, материалы: слайды "
    "лекций и семинаров я выкладываю туда по мере проведения занятий, так что всё собрано в "
    "одном месте. В-третьих, вопросы: если что-то непонятно между парами — спрашивайте в "
    "чате, не обязательно копить до семинара.\n\n"
    "Пожалуйста, вступите сегодня — так вы точно не пропустите первое объявление."
)


# ============================================================
# Orchestrate
# ============================================================

fix_slide2()
fix_slide5()
fix_slide6()
fix_failures()
new_reveal_ids = build_reveal_slides()   # appended after everything
max_id = build_max_slide()               # appended last

# --- reorder sldIdLst ---
# current order indices: 0..17 = s01..s18(base), 18 = s19(memo), 19 = s20(closing),
#                        20..25 = reveals, 26 = MAX
lst = sldIdLst()
ids = list(lst)
base_quiz_id = ids[17]
memo_id = ids[18]                # s19 памятка «AI reality-check» (итоги) — RESTORED
closing_id = ids[19]
head = ids[0:18]                 # s01..s18 base (keep order)
# desired: [s01..s04, MAX, s05..s18base, reveals..., итоги-памятка, s20 closing]
# head currently = s01,s02,s03,s04,s05,...,s18  -> insert MAX after index 3 (s04);
# memo (итоги) restored right after the reveal sequence, before the hero closing.
desired = head[0:4] + [max_id] + head[4:18] + new_reveal_ids + [memo_id] + [closing_id]
for el in list(lst):
    lst.remove(el)
for el in desired:
    lst.append(el)

prs.save(str(PPTX))
print(f"Saved {PPTX} — {len(prs.slides.__iter__.__self__._sldIdLst)} slide ids; "
      f"total slides now {len(list(prs.slides))}")
