"""
build_cheatsheets.py — 4 standalone опорные карточки Лекции 17 (capstone).

«Что студент заберёт с собой» — главный практический артефакт курса.

Output (each its own PPTX → libreoffice → PDF):
  cheatsheet-1-decision-matrix.pptx/pdf   A4 portrait  (8.27 × 11.69")
  cheatsheet-2-autonomy-ladder.pptx/pdf   A4 portrait
  cheatsheet-3-failure-modes.pptx/pdf     A4 portrait
  cheatsheet-4-master-map.pptx/pdf        A1 landscape (33.11 × 23.39")
  cheatsheets-all.pdf                      4-в-1 (pdfunite)

Source-of-truth: chapter-part4.md §5.1–§5.4 + slides/s35–s38 (русифицированные
таблицы). A1 scatter — render_master_poster.py (IDENTICAL coords из
rendered/scatter_coords.py, zero drift с deck s38).

Palette LOCKED v3: Ocean (#21295C/#065A82/#1C7293) + Teal (#028090) + Gold (#F0AB00).

MCP limitation [#55-1]: PowerPoint MCP create_presentation = 4:3 only, нет slide-size
опции → строим напрямую через python-pptx с явными slide_width/slide_height.
MCP limitation [#54-render-1]: LibreOffice добавляет drop-shadow → disable_shadow().
"""
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree
from PIL import Image

# === Palette (LOCKED v3) ===
DEEP = RGBColor(0x21, 0x29, 0x5C)
MID = RGBColor(0x06, 0x5A, 0x82)
LIGHT = RGBColor(0x1C, 0x72, 0x93)
TEAL = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF0, 0xAB, 0x00)
SLATE = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
GREEN_TINT = RGBColor(0xE4, 0xF1, 0xE8)
BLUE_TINT = RGBColor(0xE6, 0xEE, 0xF5)
SOFT_GREY = RGBColor(0xEC, 0xEF, 0xF3)
RED_FAIL = RGBColor(0xB5, 0x3A, 0x2A)   # used ONLY for ✗ glyph (semantic verdict, not palette accent)
GREEN_OK = RGBColor(0x1E, 0x7A, 0x4B)   # used ONLY for ✓ glyph

FONT = "Arial"
HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"

OK = "✓"
WARN = "⚠"
NO = "✗"

# A4 portrait / A1 landscape in inches
A4_W, A4_H = 8.268, 11.693
A1_W, A1_H = 33.110, 23.386


# --------------------------------------------------------------- helpers
def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"
    for el in sppr.findall(ns):
        sppr.remove(el)
    etree.SubElement(sppr, ns)


def new_deck(w, h):
    prs = Presentation()
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    return prs


def blank_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def text_box(slide, x, y, w, h, text, *, size=12, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, line_spacing=1.1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except Exception: pass
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def runs_box(slide, x, y, w, h, segments, *, size=12, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.1):
    """One paragraph, multiple coloured runs. segments=[(text, {opts})...]."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    try: p.line_spacing = line_spacing
    except Exception: pass
    for seg, o in segments:
        r = p.add_run(); r.text = seg
        r.font.name = o.get("font", font)
        r.font.size = Pt(o.get("size", size))
        r.font.bold = o.get("bold", False)
        r.font.italic = o.get("italic", False)
        r.font.color.rgb = o.get("color", DEEP)
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5,
                radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None, stroke_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_w)
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, *, fill=MID, stroke=None, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_w)
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def cell_text(slide, x, y, w, h, segments, *, size=10.5, align=PP_ALIGN.LEFT,
              pad=0.06, line_spacing=1.04):
    """Multi-line cell: segments is a list of lines; each line is a string or
    list of runs [(text,{opts})...]."""
    tb = slide.shapes.add_textbox(Inches(x + pad), Inches(y), Inches(w - 2 * pad), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(segments):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except Exception: pass
        runs = line if isinstance(line, list) else [(line, {})]
        for seg, o in runs:
            r = p.add_run(); r.text = seg
            r.font.name = o.get("font", FONT)
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", False)
            r.font.italic = o.get("italic", False)
            r.font.color.rgb = o.get("color", DEEP)
    return tb


def add_image_aspect(slide, path, x, y, w, h):
    p = Path(path)
    if not p.exists():
        rectangle(slide, x, y, w, h, fill=SOFT_GREY)
        text_box(slide, x, y + h / 2, w, 0.6, f"[нет: {p.name}]",
                 size=12, color=SLATE, align=PP_ALIGN.CENTER)
        return None
    with Image.open(p) as img:
        iw, ih = img.size
    img_ratio = iw / ih
    box_ratio = w / h
    if img_ratio > box_ratio:
        nw = w; nh = w / img_ratio
        cx = x; cy = y + (h - nh) / 2
    else:
        nh = h; nw = h * img_ratio
        cx = x + (w - nw) / 2; cy = y
    return slide.shapes.add_picture(str(p), Inches(cx), Inches(cy),
                                    width=Inches(nw), height=Inches(nh))


# ---- verdict-glyph run helper (✓ green / ⚠ gold / ✗ red) -------------
def verdict_runs(spec):
    """Parse a string with ✓/⚠/✗ glyphs into coloured runs."""
    out = []
    cur = ""
    for ch in spec:
        if ch in (OK, WARN, NO):
            if cur:
                out.append((cur, {"size": 10, "color": DEEP}))
                cur = ""
            col = {OK: GREEN_OK, WARN: GOLD, NO: RED_FAIL}[ch]
            out.append((ch, {"size": 12, "bold": True, "color": col}))
        else:
            cur += ch
    if cur:
        out.append((cur, {"size": 10, "color": DEEP}))
    return out


# =====================================================================
#  Shared card header / footer
# =====================================================================
def card_header(slide, w, title, subtitle, *, accent=MID):
    """Top band: accent bar + title + subtitle. Returns y for content start."""
    rectangle(slide, 0, 0, w, 0.10, fill=accent)
    text_box(slide, 0.45, 0.30, w - 0.9, 0.55, title,
             size=22, bold=True, color=DEEP)
    text_box(slide, 0.45, 0.88, w - 0.9, 0.32, subtitle,
             size=12, italic=True, color=LIGHT)
    return 1.40


def card_footer_band(slide, w, h, lines, *, fill=GOLD_TINT, stroke=GOLD,
                     top=None, height=None):
    """Bottom callout band (footer rule). lines = list of run-lines."""
    bh = height or 1.15
    by = top if top is not None else (h - bh - 0.40)
    rounded_box(slide, 0.45, by, w - 0.9, bh, fill=fill, stroke=stroke, stroke_w=1.75)
    cell_text(slide, 0.65, by + 0.10, w - 1.3, bh - 0.20, lines, size=11,
              align=PP_ALIGN.LEFT, pad=0.0, line_spacing=1.2)


def provenance_footer(slide, w, h, text):
    text_box(slide, 0.45, h - 0.34, w - 0.9, 0.26, text,
             size=9, italic=True, color=SLATE)


# =====================================================================
#  CARD 1 — Decision matrix (7 criteria) — A4 portrait
# =====================================================================
def build_card1():
    prs = new_deck(A4_W, A4_H)
    s = blank_slide(prs)
    w, h = A4_W, A4_H
    y0 = card_header(
        s, w,
        "Применять ли ИИ? — 7 критериев",
        "Опорная карточка #1 · матрица решения · итоговая лекция курса AI-usage-lessons v1.0",
        accent=MID)

    # Table: 7 rows × 4 cols. Header + 7 data rows.
    rows = [
        ("№", "Вопрос", "Индикатор-вердикт", "Пример из курса"),
        ("1", "Среда контролируемая или закрытая?",
         f"{OK} закрытая · {WARN} полу · {NO} открытая",
         "See & Spray ✓ · Monarch ✗"),
        ("2", "Данных достаточно и совпадают с эксплуатацией?",
         f"{OK} · {WARN} · {NO}",
         "компилятор-SE ✓ · Epic Sepsis ✗"),
        ("3", "Задача повторяема и высокий объём?",
         f"{OK} · {NO}",
         "Copilot ✓ · штучная архитектура ✗"),
        ("4", "Цена ошибки приемлема для ИИ?",
         f"{OK} низкая · {WARN} человек-в-петле · {NO} катастрофа",
         "Stripe Radar ✓ · CrowdStrike ✗"),
        ("5", "Эталонный отклик быстрый?",
         f"{OK} · {NO}",
         "компилятор ✓ · iBuying ✗"),
        ("6", "Нужна объяснимость?",
         f"{OK} если SHAP/LIME · {WARN} прозрачная модель · {NO} под мандат",
         "Aidoc (прозрачная) · Apple Card 2019 ✗"),
        ("7", "ИИ окупается vs базовая альтернатива?",
         f"{OK} · {NO}",
         "LaserWeeder ✓ · ML vs MPC ✗ часто"),
    ]
    # Column layout
    mx = 0.45
    cw_no = 0.45
    cw_q = 3.05
    cw_v = 2.10
    cw_e = w - 2 * mx - cw_no - cw_q - cw_v
    xs = [mx, mx + cw_no, mx + cw_no + cw_q, mx + cw_no + cw_q + cw_v]
    cws = [cw_no, cw_q, cw_v, cw_e]

    header_h = 0.50
    row_h = 1.04
    ty = y0

    # header row
    rectangle(s, mx, ty, w - 2 * mx, header_h, fill=DEEP)
    for ci, label in enumerate(rows[0]):
        text_box(s, xs[ci], ty, cws[ci], header_h, label,
                 size=11.5, bold=True, color=WHITE, align=(PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT),
                 anchor=MSO_ANCHOR.MIDDLE)
    ty += header_h

    for ri, row in enumerate(rows[1:]):
        fill = SURFACE if ri % 2 == 0 else WHITE
        rectangle(s, mx, ty, w - 2 * mx, row_h, fill=fill,
                  stroke=RGBColor(0xDD, 0xE4, 0xEC), stroke_w=0.75)
        # № cell — circle badge
        circle(s, xs[0] + cw_no / 2 - 0.16, ty + row_h / 2 - 0.16, 0.32, fill=MID)
        text_box(s, xs[0], ty, cw_no, row_h, row[0], size=13, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # question
        cell_text(s, xs[1], ty, cw_q, row_h, [row[1]], size=11.5, align=PP_ALIGN.LEFT)
        # verdict (coloured glyph runs)
        cell_text(s, xs[2], ty, cw_v, row_h, [verdict_runs(row[2])], size=10, align=PP_ALIGN.LEFT)
        # example (coloured glyph runs)
        cell_text(s, xs[3], ty, cw_e, row_h, [verdict_runs(row[3])], size=10, align=PP_ALIGN.LEFT)
        ty += row_h

    # Footer rule (GOLD highlight on the main decision rule)
    card_footer_band(
        s, w, h,
        [
            [("Правило: ", {"size": 12, "bold": True, "color": DEEP}),
             ("проходите 7 строк по порядку.", {"size": 12, "color": DEEP})],
            [("Один ", {"size": 11.5, "color": DEEP}),
             (NO, {"size": 13, "bold": True, "color": RED_FAIL}),
             (" — STOP, отказ от полного ИИ.   ", {"size": 11.5, "bold": True, "color": DEEP}),
             ("≥2 ", {"size": 11.5, "color": DEEP}),
             (WARN, {"size": 13, "bold": True, "color": GOLD}),
             (" — STOP, обоснуйте человека-в-петле + канарейку + откат.", {"size": 11.5, "color": DEEP})],
            [("Все 7 ", {"size": 11.5, "color": DEEP}),
             (OK, {"size": 13, "bold": True, "color": GREEN_OK}),
             (" — пилот с явными воротами GO/NO-GO.", {"size": 11.5, "bold": True, "color": DEEP})],
        ],
        fill=GOLD_TINT, stroke=GOLD, height=1.30, top=ty + 0.30)

    provenance_footer(
        s, w, h,
        "SHAP/LIME, HITL, MPC, GO/NO-GO — термины из материала лекции; кейсы — бренды. "
        "Источник: глава §5.1. Лекция 17 · МГТУ ИУ6.")
    return prs


# =====================================================================
#  CARD 2 — Autonomy ladder L0→L5 — A4 portrait
# =====================================================================
def build_card2():
    prs = new_deck(A4_W, A4_H)
    s = blank_slide(prs)
    w, h = A4_W, A4_H
    y0 = card_header(
        s, w,
        "Лестница автономии: L0 → L5",
        "Опорная карточка #2 · уровни автономии и критерии подъёма · итоговая лекция v1.0",
        accent=TEAL)

    rows = [
        ("Ур.", "Название", "Что ИИ делает", "Кто решает", "Критерий подъёма"),
        ("L0", "Без автоматизации", "Нет ИИ", "Человек",
         "Базовая линия собрана"),
        ("L1", "Советует", "Классиф. / предсказ. / рекоменд.", "Человек всегда",
         "Базовая линия + контроль изменений + откат"),
        ("L2", "С подтверждением", "Действует, человек подтверждает", "Человек каждое",
         "Частота ложных срабат. + канарейка + откат"),
        ("L3", "Условный (узкий домен)", "Действует в узком домене (ODD)", "На петле (HOOL)",
         "Домен (ODD) формально + телеметрия + go/no-go"),
        ("L4", "Высокий (широкий домен)", "Действует в широком домене", "Вне петли (HOTL)",
         "Надёжность 99,9% + страховка + допуск регулятора"),
        ("L5", "Полный", "Решает везде", "(в 2026 недостижим)",
         "Для большинства отраслей недоступен"),
    ]
    mx = 0.45
    cw_l = 0.62
    cw_n = 1.85
    cw_d = 1.95
    cw_w = 1.30
    cw_c = w - 2 * mx - cw_l - cw_n - cw_d - cw_w
    xs = [mx, mx + cw_l, mx + cw_l + cw_n, mx + cw_l + cw_n + cw_d,
          mx + cw_l + cw_n + cw_d + cw_w]
    cws = [cw_l, cw_n, cw_d, cw_w, cw_c]

    header_h = 0.52
    row_h = 1.06
    ty = y0
    rectangle(s, mx, ty, w - 2 * mx, header_h, fill=DEEP)
    for ci, label in enumerate(rows[0]):
        text_box(s, xs[ci] + 0.04, ty, cws[ci] - 0.08, header_h, label,
                 size=10.5, bold=True, color=WHITE,
                 align=(PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT),
                 anchor=MSO_ANCHOR.MIDDLE)
    ty += header_h

    # level badge colors: L0 grey → L5 gold (gradient of autonomy)
    badge_colors = [SLATE, MID, MID, LIGHT, TEAL, GOLD]
    for ri, row in enumerate(rows[1:]):
        fill = SURFACE if ri % 2 == 0 else WHITE
        rectangle(s, mx, ty, w - 2 * mx, row_h, fill=fill,
                  stroke=RGBColor(0xDD, 0xE4, 0xEC), stroke_w=0.75)
        # level badge
        rounded_box(s, xs[0] + 0.05, ty + row_h / 2 - 0.21, cw_l - 0.14, 0.42,
                    fill=badge_colors[ri], stroke=None, radius=0.4)
        text_box(s, xs[0], ty, cw_l, row_h, row[0], size=13, bold=True,
                 color=(DEEP if ri == 0 else WHITE) if ri != 0 else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cell_text(s, xs[1], ty, cw_n, row_h, [[(row[1], {"size": 11, "bold": True, "color": DEEP})]], align=PP_ALIGN.LEFT)
        cell_text(s, xs[2], ty, cw_d, row_h, [row[2]], size=10, align=PP_ALIGN.LEFT)
        cell_text(s, xs[3], ty, cw_w, row_h, [row[3]], size=10, align=PP_ALIGN.LEFT)
        # promotion criterion — gold for L5 "недоступен"
        crit_color = GOLD if ri == 5 else MID
        cell_text(s, xs[4], ty, cw_c, row_h,
                  [[(row[4], {"size": 9.5, "color": crit_color, "bold": (ri == 5)})]],
                  align=PP_ALIGN.LEFT)
        ty += row_h

    # Antipatterns footer band (per level)
    card_footer_band(
        s, w, h,
        [
            [("Антипаттерны: ", {"size": 11.5, "bold": True, "color": DEEP}),
             ("L1 превышение роли (Klarna) · L2 скучный надзор (Uber Tempe) · L3 расширение домена (Cruise)",
              {"size": 10.5, "color": DEEP})],
            [("L4 действие без канарейки (CrowdStrike) · L5 этический блок (LAWS) · сквозной — пропуск ступени.",
              {"size": 10.5, "color": DEEP})],
            [("Правило: ", {"size": 12, "bold": True, "color": DEEP}),
             ("определите ", {"size": 11.5, "color": DEEP}),
             ("текущий + максимально допустимый уровень", {"size": 11.5, "bold": True, "color": GOLD}),
             ("; различаются — нужен план подъёма. Большинство 2026 — L1–L2.",
              {"size": 11.5, "color": DEEP})],
        ],
        fill=TEAL_TINT, stroke=TEAL, height=1.35, top=ty + 0.30)

    provenance_footer(
        s, w, h,
        "ODD (operational design domain — рабочий домен), HOOL/HOTL (человек на/вне петли), "
        "LAWS, HITL — термины из материала. Источник: глава §5.2. Лекция 17 · МГТУ ИУ6.")
    return prs


# =====================================================================
#  CARD 3 — Failure-modes & antidotes (top-12) — A4 portrait
#  ГЛАВНАЯ карточка по практической ценности.
# =====================================================================
def build_card3():
    prs = new_deck(A4_W, A4_H)
    s = blank_slide(prs)
    w, h = A4_W, A4_H
    # Header with GOLD accent (главная карточка)
    rectangle(s, 0, 0, w, 0.10, fill=GOLD)
    text_box(s, 0.45, 0.26, w - 2.3, 0.50, "12 провалов и противоядия",
             size=22, bold=True, color=DEEP)
    # "главная карточка" badge
    rounded_box(s, w - 2.15, 0.28, 1.70, 0.42, fill=GOLD, stroke=None, radius=0.3)
    text_box(s, w - 2.15, 0.28, 1.70, 0.42, "ГЛАВНАЯ КАРТОЧКА",
             size=10, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.45, 0.80, w - 0.9, 0.30,
             "Опорная карточка #3 · реестр провалов и альтернатив · итоговая лекция v1.0",
             size=12, italic=True, color=LIGHT)
    y0 = 1.28

    # 12 rows × 4 cols (провал / источник / урок / альтернатива)
    rows = [
        ("#", "Провал", "Источник", "Урок одной фразой", "Альтернатива"),
        ("1", "Открытый мир без закрытой петли", "Zillow · Monarch · Cruise",
         "Сдвиг распределения убивает ML", "Сузить домен · человек-советник · не-ИИ"),
        ("2", "Накопление ненадёжности (агент)", "$4 200-петля · агент-SE",
         "p^N → 0 при N > 10 шагов", "Лимит бюджета + макс. шагов + точки HITL"),
        ("3", "Демо ≠ промышленная эксплуатация", "Devin · Watson · Epic · Klarna",
         "Бенчмарк вендора ≠ ваша среда", "Повторить замер на ваших данных"),
        ("4", "Скучный человек-в-петле", "Uber Tempe · F-35 ALIS",
         "Монотонный надзор проваливается", "Надзор с алертом · снизить ложн. срабат."),
        ("5", "Переавтоматизация вариативных зон", "Tesla 2018 · Boeing MAX 9",
         "Парадокс автоматизации (Bainbridge 1983)", "Jidoka: усиление, не замена"),
        ("6", "Действие без канарейки и отката", "CrowdStrike · Cloudflare",
         "Широкий домен = большой радиус поражения", "Канарейка + телеметрия + откат в 1 клик"),
        ("7", "Научная галлюцинация", "Meta Galactica",
         "Текст ≠ эксперимент", "RAG-заземление + рецензирование"),
        ("8", "Голос / видео-дипфейк", "Wendy's · Air Canada · Arup $25М",
         "Шум + сложность = провал; видео — новый вектор", "Меню по правилам · C2PA · независимый канал"),
        ("9", "Утечка обучающих данных", "Getty v. Stability · NYT v. OpenAI",
         "Хвост запоминания в моделях", "Лицензированные датасеты · аудит происхождения"),
        ("10", "Привязка к вендору (в регулируемых)", "Climate FieldView · ALIS · Watson",
         "Привязка → стратегический риск", "Гос-стиль владения · экспорт данных в договоре"),
        ("11", "Slopsquatting (цепочка поставок)", "имена npm / pip",
         "Выдуманные ИИ имена = новый вектор атаки", "Проверка SBOM + белый список импортов"),
        ("12", "Болото пилотов (90–95% не доходят)", "MIT 95% · McKinsey 5,5% · РФ 9 из 10",
         "Слип в бесконечный пилот", "Явные ворота GO/NO-GO + бюджет + базовая линия"),
    ]
    mx = 0.40
    cw_n = 0.34
    cw_f = 1.95
    cw_src = 1.70
    cw_les = 1.95
    cw_alt = w - 2 * mx - cw_n - cw_f - cw_src - cw_les
    xs = [mx, mx + cw_n, mx + cw_n + cw_f, mx + cw_n + cw_f + cw_src,
          mx + cw_n + cw_f + cw_src + cw_les]
    cws = [cw_n, cw_f, cw_src, cw_les, cw_alt]

    header_h = 0.46
    row_h = 0.69
    ty = y0
    rectangle(s, mx, ty, w - 2 * mx, header_h, fill=DEEP)
    for ci, label in enumerate(rows[0]):
        text_box(s, xs[ci] + 0.04, ty, cws[ci] - 0.08, header_h, label,
                 size=10, bold=True, color=WHITE,
                 align=(PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT),
                 anchor=MSO_ANCHOR.MIDDLE)
    ty += header_h

    for ri, row in enumerate(rows[1:]):
        fill = SURFACE if ri % 2 == 0 else WHITE
        rectangle(s, mx, ty, w - 2 * mx, row_h, fill=fill,
                  stroke=RGBColor(0xDD, 0xE4, 0xEC), stroke_w=0.6)
        # # badge
        text_box(s, xs[0], ty, cw_n, row_h, row[0], size=11, bold=True,
                 color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cell_text(s, xs[1], ty, cw_f, row_h, [[(row[1], {"size": 9, "bold": True, "color": DEEP})]], align=PP_ALIGN.LEFT, line_spacing=1.0)
        cell_text(s, xs[2], ty, cw_src, row_h, [[(row[2], {"size": 8, "italic": True, "color": SLATE})]], align=PP_ALIGN.LEFT, line_spacing=1.0)
        cell_text(s, xs[3], ty, cw_les, row_h, [[(row[3], {"size": 8.5, "color": MID})]], align=PP_ALIGN.LEFT, line_spacing=1.0)
        cell_text(s, xs[4], ty, cw_alt, row_h, [[(row[4], {"size": 8.5, "color": TEAL})]], align=PP_ALIGN.LEFT, line_spacing=1.0)
        ty += row_h

    # Footer
    card_footer_band(
        s, w, h,
        [
            [("Правило: ", {"size": 12, "bold": True, "color": DEEP}),
             ("на любой вендор-питч / планёрку / рецензию — пройдите 12 строк.",
              {"size": 11, "color": DEEP})],
            [("Узнаёте паттерн — задайте уточняющий вопрос. ", {"size": 11, "color": DEEP}),
             ("Часто этого достаточно, чтобы спасти проект.", {"size": 11, "bold": True, "color": GOLD})],
            [("12 — за пределами рабочей памяти (Miller, 7±2): держите на бумаге, не в голове.",
              {"size": 10, "italic": True, "color": SLATE})],
        ],
        fill=GOLD_TINT, stroke=GOLD, height=1.10, top=ty + 0.16)

    provenance_footer(
        s, w, h,
        "ML, HITL, RAG, SBOM, C2PA, slopsquatting (выдуманные ИИ имена-зависимости) — термины из "
        "материала; кейсы — бренды. Источник: глава §5.3. Лекция 17 · МГТУ ИУ6.")
    return prs


# =====================================================================
#  CARD 4 — Master map (16 industries) — A1 landscape poster
# =====================================================================
def build_card4():
    prs = new_deck(A1_W, A1_H)
    s = blank_slide(prs)
    w, h = A1_W, A1_H
    # Full-bleed master poster PNG (identical coords from rendered/scatter_coords.py)
    poster = ASSETS / "master-poster-a1.png"
    # poster image already contains title + axes + zones + callouts + footer.
    # Fill the whole A1 canvas with small margin.
    margin = 0.30
    add_image_aspect(s, poster, margin, margin, w - 2 * margin, h - 2 * margin)
    return prs


# =====================================================================
#  Build + export
# =====================================================================
def export_pdf(pptx_path):
    pdf_path = pptx_path.with_suffix(".pdf")
    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(pptx_path.parent), str(pptx_path)],
        check=True, capture_output=True)
    return pdf_path


def main():
    cards = [
        ("cheatsheet-1-decision-matrix", build_card1),
        ("cheatsheet-2-autonomy-ladder", build_card2),
        ("cheatsheet-3-failure-modes", build_card3),
        ("cheatsheet-4-master-map", build_card4),
    ]
    pdfs = []
    for name, fn in cards:
        prs = fn()
        pptx_path = HERE / f"{name}.pptx"
        prs.save(str(pptx_path))
        pdf = export_pdf(pptx_path)
        pdfs.append(pdf)
        print(f"Built {name}: {pptx_path.name} → {pdf.name}")

    # Combined 4-in-1
    combined = HERE / "cheatsheets-all.pdf"
    subprocess.run(["pdfunite", *[str(p) for p in pdfs], str(combined)], check=True)
    print(f"Combined: {combined.name}")


if __name__ == "__main__":
    main()
