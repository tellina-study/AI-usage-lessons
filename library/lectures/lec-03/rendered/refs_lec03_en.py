"""
Reference / page-number system for Lecture 3 — EN (issue #172, ported from
the RU refs_lec03.py; anchors re-keyed to the EN builder's visible text).

Provides:
  • URLS            — canonical URL map (ONLY from reference-registry.md).
  • SLIDE_REFS      — per display-slide source registry:
                        (num, short_name, urlkey|None, gloss[, volatile]).
  • ANCHORS         — per display-slide list of (ref_nums, anchor_substr):
                        where to inject the small superscript [N] marker(s)
                        inside the already-built visible body. anchor_substr
                        is a verbatim fragment of an EXISTING run — nothing in
                        the visible copy changes except the appended [N].
  • shrink_refs_in_frame  — post-hoc split of [N] markers into small (~52%)
                            superscript muted runs (#170-3).
  • inject_ref_markers    — walk a slide, append [N] at each ANCHORS anchor,
                            then shrink.
  • ref_list / refs_of_slide  — bottom clickable numbered source list.
  • notes_sources_block / notes_with_sources — "Sources:" block for notes,
                            [VFY]/[VFY-day-of] on volatile/unconfirmed sources.
  • page_number     — muted «N / 40» stamp bottom-right.

[VFY] policy (task req 4): course-internal / illustrative / future-dated
sources (agent-harness-registry; arXiv:2605.29420 / 2605.29463 / 2603.22489 /
2601.06007 / 2601.18699) are NEVER given as canonical primary URLs on the
visible slide — the confirmed anchor (Zheng, Gloaguen, Willison/Docker/Unit42,
Luo …) carries the [N]. The unconfirmed ones are named ONLY in notes with a
[VFY] tag. Volatile-but-real numbers get [VFY-day-of] in notes.
"""
import re

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# palette (mirror build_v3.py)
DEEP  = RGBColor(0x21, 0x29, 0x5C)
MID   = RGBColor(0x06, 0x5A, 0x82)
LIGHT = RGBColor(0x1C, 0x72, 0x93)
SLATE = RGBColor(0x5B, 0x66, 0x78)
FONT_BODY = "Arial"

# ============================================================
# URLS — canonical, ONLY from notes/research/lecture-3/reference-registry.md
# ============================================================
URLS = {
    # Air Canada / legal
    "aircanada_mccarthy": "https://www.mccarthy.ca/en/insights/blogs/techlex/moffatt-v-air-canada-misrepresentation-ai-chatbot",
    "aircanada_aba": "https://www.americanbar.org/groups/business_law/resources/business-law-today/2024-february/bc-tribunal-confirms-companies-remain-liable-information-provided-ai-chatbot/",
    # Anthropic engineering / research
    "anthropic_agents": "https://www.anthropic.com/research/building-effective-agents",
    "anthropic_context_eng": "https://www.anthropic.com/engineering/effective-context-engineering-for-ai-agents",
    "anthropic_cot_faith": "https://www.anthropic.com/research/reasoning-models-dont-say-think",
    "anthropic_retention": "https://platform.claude.com/docs/en/build-with-claude/api-and-data-retention",
    "anthropic_mcp_donate": "https://www.anthropic.com/news/donating-the-model-context-protocol-and-establishing-of-the-agentic-ai-foundation",
    # Prompt / roles / CoT
    "zheng_personas": "https://arxiv.org/abs/2311.10054",
    "wei_cot": "https://arxiv.org/abs/2201.11903",
    # Context
    "chroma_context_rot": "https://research.trychroma.com/context-rot",
    "liu_lost_middle": "https://arxiv.org/abs/2307.03172",
    # RAG
    "lewis_rag": "https://arxiv.org/abs/2005.11401",
    "barnett_7fail": "https://arxiv.org/abs/2401.05856",
    "kore_7rag": "https://www.kore.ai/blog/seven-rag-engineering-failure-points",
    "uniah": "https://arxiv.org/abs/2503.00353",
    "redhat_rag_ft": "https://www.redhat.com/en/topics/ai/rag-vs-fine-tuning",
    "ibm_rag_ft": "https://www.ibm.com/think/topics/rag-vs-fine-tuning",
    "bigdata_ft": "https://bigdataboutique.com/blog/fine-tuning-llms-when-rag-isnt-enough",
    # Fine-tuning / PEFT / distillation
    "hf_beyond_lora": "https://huggingface.co/blog/peft-beyond-lora",
    "lora_paper": "https://arxiv.org/abs/2106.09685",
    "qlora_paper": "https://arxiv.org/abs/2305.14314",
    "hinton_distill": "https://arxiv.org/abs/1503.02531",
    "peft_survey": "https://arxiv.org/abs/2403.14608",
    "luo_forgetting": "https://arxiv.org/abs/2308.08747",
    # Agents / loop
    "yao_react": "https://arxiv.org/abs/2210.03629",
    "cognition_no_multiagent": "https://cognition.ai/blog/dont-build-multi-agents",
    "gloaguen_presence": "https://arxiv.org/abs/2602.11988",
    "claude_code_51735": "https://github.com/anthropics/claude-code/issues/51735",
    # Security / retention
    "willison_mcp_inject": "https://simonwillison.net/2025/Apr/9/mcp-prompt-injection/",
    "docker_mcp_horror": "https://www.docker.com/blog/mcp-horror-stories-github-prompt-injection/",
    "unit42_mcp": "https://unit42.paloaltonetworks.com/model-context-protocol-attack-vectors/",
    "nyt_openai_bloomberg": "https://news.bloomberglaw.com/ip-law/openai-must-turn-over-20-million-chatgpt-logs-judge-affirms",
    # Failures / postmortems / market
    "jain_4200": "https://medium.com/@sattyamjain96/the-agent-that-burned-4-200-in-63-hours-a-production-ai-postmortem-d38fd9586a85",
    "mindstudio_reliability": "https://www.mindstudio.ai/blog/reliability-compounding-problem-ai-agent-stacks",
    "mit_nanda_fortune": "https://fortune.com/2025/08/18/mit-report-95-percent-generative-ai-pilots-at-companies-failing-cfo/",
}

# ============================================================
# SLIDE_REFS — (num, short_name, urlkey|None, gloss[, volatile])
# volatile=True → [VFY-day-of] in notes. urlkey=None → course-internal
# (agent-harness-registry, unconfirmed) → notes-only [VFY], never on slide.
# ============================================================
SLIDE_REFS = {
    "s01": [
        ("1", "McCarthy Tétrault — Moffatt v. Air Canada (BC CRT, 14.02.2024)", "aircanada_mccarthy",
         "the bot invented a refund policy; tribunal: the company is liable for the bot's answer"),
        ("2", "ABA Business Law Today — companies are liable for their AI chatbot", "aircanada_aba",
         "\"the bot is not a separate legal entity\": wrong architecture choice for a lookup task"),
    ],
    "s05": [
        ("1", "Anthropic — Building Effective Agents (\"find the simplest\")", "anthropic_agents",
         "don't complicate the architecture beyond the task's requirement — burden-of-proof allocation"),
    ],
    "s05a": [
        ("1", "Zheng et al. 2024, Findings of EMNLP — personas don't improve factual accuracy", "zheng_personas",
         "162 personas · 8 domains · 2410 factual questions · 4 families → no accuracy gain"),
    ],
    "s06": [
        ("1", "Wei et al. 2022 — Chain-of-Thought Prompting", "wei_cot",
         "step-by-step reasoning raises reliability on arithmetic / multi-step logic"),
        ("2", "Anthropic — Reasoning Models Don't Always Say What They Think", "anthropic_cot_faith",
         "faithfulness: Claude 3.7 ~25%, DeepSeek R1 ~39% — the reasoning need not reflect the real cause", True),
    ],
    "s08": [
        ("1", "Chroma Research — Context Rot", "chroma_context_rot",
         "retrieval accuracy drops as the number of tokens in context grows"),
        ("2", "Liu et al. 2023 — Lost in the Middle", "liu_lost_middle",
         "the same \"lost in the middle\" phenomenon from Lecture 2 — a new term, not a new entity"),
        ("3", "Anthropic — Effective Context Engineering", "anthropic_context_eng",
         "minimal high-signal context is an engineering requirement, not aesthetics"),
    ],
    "s10": [
        ("1", "Lewis et al. 2020 — Retrieval-Augmented Generation (NeurIPS)", "lewis_rag",
         "the canonical RAG paper: indexing → retrieval → grounded generation"),
    ],
    "s11": [
        ("1", "IBM — RAG vs Fine-Tuning (vendor-neutral)", "ibm_rag_ft",
         "signals-for-RAG: large / changing / provenance / private"),
        ("2", "U-NIAH — RAG win-rate is higher for smaller models", "uniah",
         "RAG's advantage over a direct answer is especially large for smaller models", True),
    ],
    "s12": [
        ("1", "Red Hat — RAG vs Fine-Tuning (when not RAG)", "redhat_rag_ft",
         "corpus fits the window → full-context+cache; fixed value → lookup; live → API without an index"),
        ("2", "McCarthy Tétrault — Air Canada (\"generation on top of a fixed policy\")", "aircanada_mccarthy",
         "a fixed policy → deterministic lookup, not retrieval+generation"),
    ],
    "s13": [
        ("1", "Barnett et al. 2024 — Seven Failure Points (RAG)", "barnett_7fail",
         "\"returned something ≠ returned the right thing\": 7 failure points of RAG engineering"),
        ("2", "Kore.ai — Seven RAG Engineering Failure Points", "kore_7rag",
         "legal-AI / medical-RAG / support bot — degradation at scale without observability"),
        ("3", "McCarthy Tétrault — Air Canada (grounding failure)", "aircanada_mccarthy",
         "generated text in a role that required a retrieved, verified fact"),
    ],
    "s13b": [
        ("1", "IBM — RAG vs Fine-Tuning (fine-tuning changes the weights)", "ibm_rag_ft",
         "prompt/RAG change the context; fine-tuning changes the model's weights themselves"),
    ],
    "s14": [
        ("1", "BigData Boutique — Fine-Tuning When RAG Isn't Enough", "bigdata_ft",
         "fine-tuning narrowed to behavior/style/format/policy; knowledge → RAG", True),
        ("2", "Hinton, Vinyals, Dean 2015 — Distilling the Knowledge", "hinton_distill",
         "distillation is a standalone technique, taxonomically NOT a kind of fine-tuning"),
        ("3", "PEFT survey — taxonomy of tuning methods", "peft_survey",
         "surveys place distillation and fine-tuning in different categories"),
    ],
    "s15": [
        ("1", "Hu et al. 2021 — LoRA", "lora_paper",
         "base weights are frozen, low-rank adapters are trained"),
        ("2", "Dettmers et al. 2023 — QLoRA", "qlora_paper",
         "LoRA on top of a quantized model → fine-tuning on a single GPU"),
        ("3", "HF PEFT — Beyond LoRA (LoRA 98.4% of 20,834 cards)", "hf_beyond_lora",
         "98.4% of PEFT-tagged models use LoRA; caveat: share among PEFT-tagged only", True),
    ],
    "s16": [
        ("1", "Luo et al. 2023 — Catastrophic Forgetting in LLM Continual FT", "luo_forgetting",
         "narrow aggressive FT breaks general abilities; worse as the model scale grows"),
    ],
    "s19": [
        ("1", "Anthropic — MCP donation / Agentic AI Foundation (N×M→N+M)", "anthropic_mcp_donate",
         "MCP standardizes connection; ease of connection ≠ safety of what you connect", True),
    ],
    "s21": [
        ("1", "Yao et al. 2022 — ReAct (plan→act→check→iterate)", "yao_react",
         "interleaving reasoning and actions; every step is a place of failure"),
        ("2", "Anthropic — Reasoning Models faithfulness (check ≠ self-assessment)", "anthropic_cot_faith",
         "the check step is validation against an external criterion, not the model's self-assessment", True),
    ],
    "s22": [
        ("1", "Anthropic — Building Effective Agents (workflow vs agent)", "anthropic_agents",
         "workflow = predefined paths; agent = dynamic process; latency/cost↔quality trade-off"),
        ("2", "Cognition — Don't Build Multi-Agents", "cognition_no_multiagent",
         "multi-agent by default is not an upgrade; fragility of parallel subagents"),
    ],
    "s22b": [
        ("1", "agent-harness-registry — map of the agent's equipment slots", None,
         "source not confirmed by canonical URL 2026-08-30; Claude Code/Cursor/Aider — vendor sites, verify day-of", True),
    ],
    "s22c": [
        ("1", "agent-harness-registry (live-eval) — the agent memory spectrum", None,
         "source not confirmed by canonical URL 2026-08-30; parallel with the RAG scale criterion", True),
    ],
    "s22d": [
        ("1", "agent-harness-registry (live-eval) — Letta Tier D / Memory Tool 17% tail", None,
         "source not confirmed by canonical URL 2026-08-30; numbers volatile (1.0/0.833/0.750; 17%)", True),
    ],
    "s22e": [
        ("1", "Gloaguen et al. 2026 — Evaluating AGENTS.md (presence paradox)", "gloaguen_presence",
         "the presence of an instruction file gives no significant gain; helps in a documentation gap", True),
        ("2", "GitHub anthropics/claude-code#51735 (the error repeated after 25 days)", "claude_code_51735",
         "a written record of a past mistake did not prevent it from recurring"),
    ],
    "s23": [
        ("1", "Sattyam Jain 2026 — The Agent That Burned $4,200 in 63 Hours", "jain_4200",
         "a loop with no limits on HTTP 429; a retry script would have solved it almost for free", True),
        ("2", "MindStudio — Reliability Compounding Problem", "mindstudio_reliability",
         "5×99%≈95%, 10→90%, 20→82% — reliabilities multiply"),
        ("3", "Cognition — Don't Build Multi-Agents (fragility)", "cognition_no_multiagent",
         "dependent subtasks → parallel subagents make conflicting decisions"),
    ],
    "s25": [
        ("1", "Docker — MCP Horror Stories: GitHub Prompt Injection", "docker_mcp_horror",
         "GitHub MCP heist: an issue instruction + a broad token → exfiltration of private repos"),
        ("2", "Simon Willison — Prompt injection via MCP", "willison_mcp_inject",
         "the model doesn't tell data from a command; untrusted content = a command"),
        ("3", "Palo Alto Unit 42 — MCP Attack Vectors", "unit42_mcp",
         "tool poisoning / every connection = a new trust boundary"),
        ("4", "Bloomberg Law — NYT v. OpenAI (court ordered logs to be kept)", "nyt_openai_bloomberg",
         "ZDR doesn't cover everything; a court order overrides any retention policy"),
        ("5", "Anthropic — API and Data Retention (the bounds of ZDR)", "anthropic_retention",
         "ZDR doesn't cover third-party / MCP connectors — exactly what an agent is built from", True),
    ],
    "s25b": [
        ("1", "agent-harness-registry — a survey through the equipment frame", None,
         "source not confirmed; Claude Code/Aider/Cursor/OpenHands — vendor sites, verify day-of", True),
    ],
    "s26": [
        ("1", "Anthropic — Building Effective Agents (the ladder rule)", "anthropic_agents",
         "stay on the bottom rung; every step up is a trade, not an improvement"),
    ],
    "s27": [
        ("1", "Anthropic — Building Effective Agents (the choice route)", "anthropic_agents",
         "a top-down route of questions; step 1 — a deterministic task → plain code, STOP"),
        ("2", "McCarthy Tétrault — Air Canada (the bottom row of the matrix)", "aircanada_mccarthy",
         "a generative architecture for a deterministic task = a bottom-row error"),
    ],
    "s29": [
        ("1", "Anthropic — Reasoning Models faithfulness (self-rationale ≠ control)", "anthropic_cot_faith",
         "a human checks the result/facts against the source, not the model's self-explanation", True),
        ("2", "MIT NANDA — State of AI in Business 2025 (~95% of pilots with no ROI)", "mit_nanda_fortune",
         "the root is the learning gap and integration failure, not model quality; a report, not a law", True),
    ],
}

# ============================================================
# ANCHORS — where [N] markers go inside the already-built visible body.
# (ref_nums:str, anchor_substr:str). anchor_substr MUST be a verbatim
# fragment of an existing run; the marker «[N]» is appended right after it.
# Nothing else in the visible copy changes.
# NB: s22c/s22d/s25b anchors carry a course-internal source (urlkey None) →
# still get a visible [N] whose bottom-list entry has NO hyperlink (framed as
# "independent live-eval registry", with [VFY] only in notes).
# ============================================================
# NB (issue #172): anchor substrings are VERBATIM fragments of the EN builder's
# rendered runs. The EN builder uses CURLY quotes/apostrophes (“ ” ‘ ’), so
# anchors are deliberately chosen quote-free to avoid U+2019/U+201C mismatch.
ANCHORS = {
    "s01": [("1,2", "the company refunds the difference")],
    "s05": [("1", "without a reason expressed in the task")],
    "s05a": [("1", "improve answer accuracy compared with answering with no persona at all")],
    "s06": [("1", "23 − 7 = 16"),
            ("2", "the share of cases where the model mentioned the hint it actually used")],
    "s08": [("1,2", "a new term, not a new phenomenon"),
            ("3", "that is an engineering requirement, not aesthetics")],
    "s10": [("1", "is a correct answer from a RAG system")],
    "s11": [("1", "a strong signal on the features below"),
            ("2", "A single feature is a reason to look closer")],
    "s12": [("1", "RAG is redundant if ANY of the three holds"),
            ("2", "A fixed policy / value")],
    "s13": [("1,2", "RAG has no"),
            ("3", "Air Canada — architecture breakdown")],
    "s13b": [("1", "change the model itself")],
    "s14": [("1", "Fine-tuning is not dead"),
            ("2,3", "Distillation is a separate technique, often applied in tandem with fine-tuning")],
    "s15": [("1,2", "low-rank adapter matrices; QLoRA — the same on top of a quantized model"),
            ("3", "of models tagged PEFT are LoRA")],
    "s16": [("1", "degradation of the model")],
    "s19": [("1", "N+M")],
    "s21": [("1", "works in a loop, defining the sequence of steps itself"),
            ("2", "validation against an EXTERNAL criterion — not the model")],
    "s22": [("1", "A predictable task → workflow"),
            ("2", "the value justifies a multiple increase → agent")],
    "s22b": [("1", "Five typical slots")],
    "s22c": [("1", "The spectrum runs from a flat file to graph databases")],
    "s22d": [("1", "sometimes it is dramatically not")],
    "s22e": [("1", "the mere presence of an instruction file gives NO significant gain"),
             ("2", "did NOT prevent it from recurring 25 days later")],
    "s23": [("1", "$4,200 over 63 hours"),
            ("2", "fewer hops + validation between steps"),
            ("3", "parallel subagents make conflicting")],
    "s25": [("1", "dumped private repositories into a public PR"),
            ("2,3", "a surface appears that a single call did not have"),
            ("4,5", "a court order (NYT v. OpenAI) + third-party/MCP outside ZDR")],
    "s25b": [("1", "which harness slots are filled and where the agent physically lives")],
    "s26": [("1", "Every climb is a TRADE-OFF")],
    "s27": [("1", "stop at the first question that triggers"),
            ("2", "if the task is deterministic and verifiable — plain code, STOP here")],
    "s29": [("1", "NOT by the plausibility of the reasoning"),
            ("2", "the root is in the learning gap and integration failure, not in model quality")],
}

# ============================================================
# NOTES_ANCHORS — where [N] markers go inside the READABLE speaker notes
# (.md «## Speaker notes»). (ref_nums:str, notes_anchor_substr:str) — the
# marker «[N]» is appended right after the first verbatim occurrence.
# ============================================================
# NB (issue #172): NOTES_ANCHORS is NOT used by the EN builder at build time —
# the EN speaker notes are translated with the [N] markers already positioned
# in-place (Sources: block baked into slides-en/*.md), so patch_notes need not
# re-run. Kept translated for consistency / possible re-patching.
NOTES_ANCHORS = {
    "s01": [("1,2", "responsible for its own actions.\"")],
    "s05": [("1", "one call to the model with a well-composed prompt")],
    "s05a": [("1", "whether this shifts factual accuracy")],
    "s06": [("1", "16 plus 12 equals 28"),
            ("2", "roughly two in five")],
    "s08": [("3", "visible to the model at inference")],
    "s10": [("1", "grounding on these fragments")],
    "s11": [("1", "a strong signal on one or more features")],
    "s12": [("1", "Three explicit \"not RAG\" criteria")],
    "s13": [("1,2", "does not mean \"the system returned the right thing\""),
            ("3", "Air Canada")],
    "s13b": [("1", "continued training of an already finished, pretrained model")],
    "s14": [("1", "it narrowed"),
            ("2,3", "distillation")],
    "s15": [("1", "adapters"),
            ("3", "98")],
    "s16": [("1", "degradation of the model's general abilities")],
    "s19": [("1", "an open standard for a single way to connect")],
    "s21": [("1", "works in a loop"),
            ("2", "the check step")],
    "s22": [("1", "separates two concepts"),
            ("2", "multi-agent")],
    "s22b": [("1", "the loop plus the harness")],
    "s22c": [("1", "a flat file")],
    "s22d": [("1", "an independent registry")],
    "s22e": [("1", "presence paradox"),
             ("2", "claude-code")],
    "s23": [("1", "429"),
            ("2", "multiply"),
            ("3", "multi-agent")],
    "s25": [("1", "private repositories and publish them here.\""),
            ("2,3", "doesn't tell data from a command"),
            ("4,5", "a court order")],
    "s25b": [("1", "the equipment frame")],
    "s26": [("1", "the ladder of architectural complexity")],
    "s27": [("1", "stopping at the first that fires"),
            ("2", "deterministic and verifiable")],
    "s29": [("1", "unfaithful"),
            ("2", "learning gap")],
}


def inject_notes_markers(note_text, sid):
    """Insert [N] markers into readable notes at NOTES_ANCHORS anchors.
    Returns (new_text, missed:list[str])."""
    missed = []
    for ref_nums, anchor in NOTES_ANCHORS.get(sid, []):
        marker = f"[{ref_nums}]"
        if marker in note_text:
            continue
        i = note_text.find(anchor)
        if i < 0:
            missed.append(anchor)
            continue
        j = i + len(anchor)
        note_text = note_text[:j] + marker + note_text[j:]
    return note_text, missed


# ============================================================
# [N] shrink (#170-3): split [N] markers into small superscript muted runs.
# ============================================================
_REF_RE = re.compile(r'\[\d+(?:\s*[,–—-]\s*\d+)*\]')
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _run_props(src_run):
    f = src_run.font
    sz = f.size
    return {
        "name": f.name,
        "size_pt": (sz.pt if sz is not None else None),
        "bold": f.bold,
        "italic": f.italic,
        "color": (f.color.rgb if (f.color and f.color.type is not None) else None),
    }


def _clone_run_after(anchor_r, props, text, *, ref=False, ref_frac=0.52,
                     ref_color=LIGHT):
    new_r = etree.SubElement(anchor_r.getparent(), _A + "r")
    anchor_r.addnext(new_r)
    rpr = etree.SubElement(new_r, _A + "rPr")
    base = props["size_pt"] or 16.0
    if ref:
        rpr.set("sz", str(int(round(base * ref_frac * 100))))
        rpr.set("baseline", "30000")
        rpr.set("b", "0")
        rpr.set("i", "1")
    else:
        if props["size_pt"] is not None:
            rpr.set("sz", str(int(round(base * 100))))
        if props["bold"] is not None:
            rpr.set("b", "1" if props["bold"] else "0")
        if props["italic"] is not None:
            rpr.set("i", "1" if props["italic"] else "0")
    if props["name"]:
        for tag in ("latin", "cs", "ea"):
            el = etree.SubElement(rpr, _A + tag)
            el.set("typeface", props["name"])
    col = ref_color if ref else props["color"]
    if col is not None:
        fill = etree.SubElement(rpr, _A + "solidFill")
        clr = etree.SubElement(fill, _A + "srgbClr")
        clr.set("val", str(col))
    t = etree.SubElement(new_r, _A + "t")
    t.text = text
    return new_r


def shrink_refs_in_frame(text_frame, *, ref_frac=0.52, ref_color=LIGHT):
    for para in text_frame.paragraphs:
        for run in list(para.runs):
            txt = run.text
            if not txt or "[" not in txt:
                continue
            matches = list(_REF_RE.finditer(txt))
            if not matches:
                continue
            props = _run_props(run)
            run.text = txt[:matches[0].start()]
            anchor = run._r
            for i, m in enumerate(matches):
                anchor = _clone_run_after(anchor, props, m.group(), ref=True,
                                          ref_frac=ref_frac, ref_color=ref_color)
                nxt = matches[i + 1].start() if i + 1 < len(matches) else len(txt)
                between = txt[m.end():nxt]
                if between:
                    anchor = _clone_run_after(anchor, props, between, ref=False)
    return text_frame


# ============================================================
# inject_ref_markers — append [N] at ANCHORS anchors on a slide, then shrink.
# Returns list of (anchor_substr, matched?) for verification.
# ============================================================
def _iter_frames(slide):
    for shp in slide.shapes:
        if shp.has_text_frame:
            yield shp.text_frame


def inject_ref_markers(slide, sid):
    anchors = ANCHORS.get(sid, [])
    report = []
    for ref_nums, anchor in anchors:
        marker = f"[{ref_nums}]"
        placed = False
        for tf in _iter_frames(slide):
            if placed:
                break
            for para in tf.paragraphs:
                if placed:
                    break
                for run in para.runs:
                    if anchor in run.text and marker not in run.text:
                        run.text = run.text.replace(anchor, anchor + marker, 1)
                        placed = True
                        break
        report.append((anchor, placed))
    # shrink every frame that now carries a [N]
    for tf in _iter_frames(slide):
        shrink_refs_in_frame(tf)
    return report


# ============================================================
# ref_list / refs_of_slide — bottom clickable numbered source list.
# ============================================================
def ref_list(slide, entries, *, y=7.06, x=0.55, w=12.25, h=0.36, size=8.5,
             color=LIGHT, line_spacing=1.02, tail=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    for i, (num, name, url) in enumerate(entries):
        rm = p.add_run(); rm.text = f"[{num}] "
        rm.font.name = FONT_BODY; rm.font.size = Pt(size)
        rm.font.bold = True; rm.font.italic = True; rm.font.color.rgb = MID
        rn = p.add_run(); rn.text = name
        rn.font.name = FONT_BODY; rn.font.size = Pt(size)
        rn.font.italic = True; rn.font.color.rgb = color
        if url:
            try:
                rn.hyperlink.address = url
            except Exception:
                pass
        if i < len(entries) - 1:
            rs = p.add_run(); rs.text = "   ·   "
            rs.font.name = FONT_BODY; rs.font.size = Pt(size)
            rs.font.italic = True; rs.font.color.rgb = color
    if tail:
        rt = p.add_run(); rt.text = "   ·   " + tail
        rt.font.name = FONT_BODY; rt.font.size = Pt(size)
        rt.font.italic = True; rt.font.color.rgb = SLATE
    return tb


def _resolve_refs(sid):
    out = []
    for entry in SLIDE_REFS.get(sid, []):
        num, name, urlkey, gloss = entry[0], entry[1], entry[2], entry[3]
        volatile = len(entry) > 4 and entry[4]
        url = URLS.get(urlkey, "") if urlkey else ""
        out.append((num, name, url, gloss, volatile))
    return out


def refs_of_slide(slide, sid, *, y=7.06, x=0.55, w=12.25, tail=None):
    resolved = _resolve_refs(sid)
    if not resolved:
        return None
    entries = [(num, name, url) for (num, name, url, gloss, vol) in resolved]
    n = len(entries)
    size = 8.5 if n <= 3 else (8.0 if n <= 4 else 7.4)
    return ref_list(slide, entries, y=y, x=x, w=w, size=size, tail=tail)


# ============================================================
# notes "Sources:" block
# ============================================================
def notes_sources_block(sid):
    resolved = _resolve_refs(sid)
    if not resolved:
        return ""
    lines = ["Sources:"]
    for (num, name, url, gloss, vol) in resolved:
        if url:
            vfy = " [VFY-day-of]" if vol else ""
            lines.append(f"[{num}] {name} — {gloss}. {url}{vfy}")
        else:
            # course-internal / unconfirmed → [VFY], no canonical URL on record
            lines.append(f"[{num}] {name} — {gloss}. [VFY: not confirmed by canonical URL, "
                         f"present as data from an independent live-eval registry, not as a primary source]")
    return "\n".join(lines)


# ============================================================
# page number — muted «N / TOTAL» bottom-right
# ============================================================
def page_number(slide, n, total=None, *, color=SLATE):
    txt = f"{n} / {total}" if total else str(n)
    tb = slide.shapes.add_textbox(Inches(12.33), Inches(7.16), Inches(0.95),
                                  Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; p.line_spacing = 1.0
    r = p.add_run(); r.text = txt
    r.font.name = FONT_BODY; r.font.size = Pt(10); r.font.italic = True
    r.font.color.rgb = color
    return tb
