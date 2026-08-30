"""
Full 43-slide EN build of Lecture 10 "AI in Agriculture".

EN render of build_lec10.py (issue #172 Ф3). Source-of-truth unchanged:
deck.yaml v1 (43 slides) + chapter v3.1 + slides-en/*.md (43 EN readable speaker notes).

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: "Ocean rounded box" (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9).

Build via: python3 build_lec10_en.py — generates lec-10-en.pptx (43 slides).
Speaker notes loaded from slides-en/*.md. Visible strings live in build_lec10_pN_en.py.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
LIGHT_TINT = RGBColor(0xE6, 0xEE, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
RED_WARN = RGBColor(0xC0, 0x39, 0x2B)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered" / "assets"
OUT = ROOT / "rendered" / "lec-10-en.pptx"
SLIDES_DIR = ROOT / "slides-en"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


# ========== Helpers ==========

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def hr_line(slide, x, y, w, color=LIGHT, weight=1.0):
    shp = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    shp.line.color.rgb = color
    shp.line.width = Pt(weight)
    return shp


def add_arrow(slide, x, y, w, h, *, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w), height=Inches(h))
    elif w:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w))
    else:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


# ========== Slide-level helpers ==========

# 7 sections: 0 Opening, 1 Field, 2 Robot, 3 Animal, 4 Chain, 4-bis Environment, 5 Shelf
SECTIONS = ["Opening", "L1 Field", "L2 Robot", "L3 Animal", "L4 Chain", "Environment", "L5 Shelf"]


def add_progress_bar(slide, current_section, *, y=6.95):
    """7-card progress bar для section dividers и cover.
    current_section=-1 → no card highlighted (e.g., cover slide)."""
    sections = SECTIONS
    n = len(sections)
    bar_x = 0.5
    bar_w = 12.33
    gap = 0.08
    card_w = (bar_w - gap * (n - 1)) / n
    card_h = 0.32
    for i, name in enumerate(sections):
        cx = bar_x + i * (card_w + gap)
        is_current = (i == current_section and current_section >= 0)
        if is_current:
            ocean_box(slide, cx, y, card_w, card_h,
                      fill=GOLD, stroke=GOLD, stroke_pt=0.0, radius_pt=6)
            text_box(slide, cx, y+0.04, card_w, card_h-0.06,
                     f"{i}. {name}", size=10, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            ocean_box(slide, cx, y, card_w, card_h,
                      fill=SURFACE, stroke=LIGHT, stroke_pt=0.5, radius_pt=6)
            text_box(slide, cx, y+0.04, card_w, card_h-0.06,
                     f"{i}. {name}", size=10, bold=False, color=LIGHT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text):
    """Standard footer для caption / source 12pt italic."""
    text_box(slide, 0.5, 7.05, 12.33, 0.30, text,
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)


def add_assertion_title(slide, text, *, size=24, y=0.4):
    """Assertion as page title — 24-26pt bold deep, left aligned."""
    text_box(slide, 0.6, y, 12.13, 1.2, text,
             size=size, bold=True, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.1)


def load_speaker_notes(slide_id):
    """Extract speaker notes from slides-en/sNN-*.md file."""
    slides_dir = SLIDES_DIR
    matches = list(slides_dir.glob(f"{slide_id}-*.md"))
    if not matches:
        return ""
    content = matches[0].read_text(encoding="utf-8")
    m = re.search(r"## Speaker notes\s*\n+(.*?)(?:\n## |\Z)", content, re.DOTALL)
    if not m:
        return ""
    notes = m.group(1).strip()
    # Strip markdown formatting
    notes = re.sub(r"\*\*([^*]+)\*\*", r"\1", notes)
    notes = re.sub(r"\*([^*]+)\*", r"\1", notes)
    return notes


def section_divider(prs, number, title, frame_phrase, current_section, caption):
    """Generic section divider slide with progress bar."""
    s = blank(prs); set_slide_bg(s, WHITE)

    # Huge decorative number on left
    text_box(s, 0.0, 0.7, 5.5, 6.0, str(number),
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD, line_spacing=0.9)

    # Right: title + frame
    text_box(s, 5.8, 2.5, 7.0, 1.5, title,
             size=34, bold=True, color=DEEP, line_spacing=1.1)
    text_box(s, 5.8, 4.4, 7.0, 1.2, frame_phrase,
             size=18, italic=True, color=MID, line_spacing=1.25)

    # Caption above progress bar
    text_box(s, 0.6, 6.55, 12.13, 0.35, caption,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)

    # Progress bar
    add_progress_bar(s, current_section=current_section, y=6.95)
    return s


# ========== Main entry — orchestrates all 43 slide builders ==========

def main():
    prs = setup_pres()

    # Phase 6 — building all 43 slides
    # Import slide builders from modular EN parts to keep each file <600 lines
    from build_lec10_p1_en import build_part1_full  # s01-s14
    from build_lec10_p2_en import build_part2  # s15-s21
    from build_lec10_p3_en import build_part3  # s22-s38

    helpers = {
        "blank": blank, "set_slide_bg": set_slide_bg, "text_box": text_box,
        "text_runs": text_runs, "ocean_box": ocean_box, "filled_rect": filled_rect,
        "hr_line": hr_line, "add_arrow": add_arrow, "add_image": add_image,
        "add_speaker_notes": add_speaker_notes, "add_progress_bar": add_progress_bar,
        "add_footer": add_footer, "add_assertion_title": add_assertion_title,
        "load_speaker_notes": load_speaker_notes, "section_divider": section_divider,
        "disable_shadow": disable_shadow,
        # Palette
        "DEEP": DEEP, "MID": MID, "LIGHT": LIGHT, "TEAL": TEAL, "SURFACE": SURFACE,
        "WHITE": WHITE, "GOLD": GOLD, "SLATE": SLATE, "COVER_OUTLINE": COVER_OUTLINE,
        "GOLD_TINT": GOLD_TINT, "TEAL_TINT": TEAL_TINT, "LIGHT_TINT": LIGHT_TINT,
        "SOFT_GREY": SOFT_GREY, "DARK_GREY": DARK_GREY, "RED_WARN": RED_WARN,
        # Constants
        "FONT_HEAD": FONT_HEAD, "FONT_BODY": FONT_BODY, "ASSETS": ASSETS,
        "SECTIONS": SECTIONS,
        # PPT
        "MSO_SHAPE": MSO_SHAPE, "MSO_ANCHOR": MSO_ANCHOR, "PP_ALIGN": PP_ALIGN,
        "Inches": Inches, "Pt": Pt, "RGBColor": RGBColor,
    }

    build_part1_full(prs, helpers)
    build_part2(prs, helpers)
    build_part3(prs, helpers)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
