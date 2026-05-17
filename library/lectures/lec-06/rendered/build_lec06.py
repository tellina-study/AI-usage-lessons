"""
Full 32-slide build of Лекции 6 «AI в инженерном проектировании и CAD/CAM»
(Phase 6 visual loop).
(Папка репо lec-06 совпадает с номером лекции по плану курса — 6.)

Source-of-truth: deck.yaml (v1, 32 слайда) + chapter v2 (status=reviewed,
~12,860 слов) + slides/*.md (32 файла с readable speaker notes 150-300 слов).

Issue #101. Phase 6 структурная правка: +6-й section-divider перед Частью 6
(s29 «Часть 6. Синтез»); worked-decision→s30, правило+матрица→s31, Q&A→s32.

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).

Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml = 75.0 мин.

Helper layer adapted from lec-07 build (proven Phase 6 pipeline).

Build via: python3 build_lec06.py — generates lec-06.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path("/tmp/lec-06-wt/library/lectures/lec-06")
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-06.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


# ============================================================
# Helpers
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            p = tf.add_paragraph()
            p.alignment = cfg.get("align", align)
            p.line_spacing = cfg.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def down_arrow(slide, x, y, w, h, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def icon(name, color="blue"):
    """Return Path to recolored icon PNG. color: blue|teal|gold|white."""
    p = ICONS / f"{name}-{color}.png"
    return p


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    if not Path(path).exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path)
            iw, ih = img.size
            img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        ir = iw / ih
        br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x), Inches(y + (h - ah) / 2),
                                     width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2), Inches(y),
                                     height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.45, h=1.05, w=12.3, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.15, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def subhead(slide, text, *, y=1.20, size=15, color=MID, x=0.55, w=12.3):
    text_box(slide, x=x, y=y, w=w, h=0.40, text=text,
             size=size, italic=True, color=color, line_spacing=1.20)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.4,
                radius=True, radius_adj=0.14)
    text_box(slide, x=x + 0.25, y=y + 0.06, w=w - 0.5, h=h - 0.12, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.22)


def footer(slide, text):
    text_box(slide, x=0.55, y=7.04, w=12.3, h=0.35, text=text,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             line_spacing=1.05)


def speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding='utf-8')
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Section divider — unified template (7 cards, Лекция 1/7 pattern)
# Лекция 6: 7 частей (0..6). roadmap-bar только на divider'ах + cover.
# ============================================================

NAV_SECTIONS = [
    ("0", "Старт\n+ карта"),
    ("1", "Фундамент\n+ таксономия"),
    ("2", "Генеративный\nдизайн ≠ genAI"),
    ("3", "Суррогат\n/ PINN"),
    ("4", "Генеративный\nAI / LLM"),
    ("5", "Когда\nсказать «нет»"),
    ("6", "Синтез"),
]


def build_section_divider(p, here_idx, title, frame_phrase, notes_slide_id):
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=0.55, y=1.0, w=4.2, h=5.0, text=str(here_idx),
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=4.85, y=1.55, w=8.0, h=0.5, text="ЧАСТЬ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    text_box(s, x=4.85, y=2.10, w=8.0, h=1.9, text=title,
             size=44, bold=True, color=DEEP, align=PP_ALIGN.LEFT,
             line_spacing=1.12)
    filled_rect(s, 4.85, 4.05, 0.05, 0.85, fill=TEAL)
    text_box(s, x=5.05, y=4.05, w=7.6, h=1.3, text=frame_phrase,
             size=18, color=MID, align=PP_ALIGN.LEFT, line_spacing=1.30)
    bar_y = 5.85
    n = 7
    total_w = 12.3
    gap = 0.12
    cw = (total_w - gap * (n - 1)) / n
    sx = 0.55
    ch = 0.95
    for i, (num, sec_title) in enumerate(NAV_SECTIONS):
        x = sx + i * (cw + gap)
        if i == here_idx:
            ocean_box(s, x, bar_y, cw, ch, fill=GOLD, stroke=GOLD, stroke_pt=2.0)
            nc, tc = WHITE, WHITE
        elif i < here_idx:
            ocean_box(s, x, bar_y, cw, ch, fill=TEAL_TINT, stroke=TEAL,
                      stroke_pt=1.2)
            nc, tc = TEAL, MID
        else:
            ocean_box(s, x, bar_y, cw, ch, fill=WHITE, stroke=LIGHT,
                      stroke_pt=1.0)
            nc, tc = LIGHT, SLATE
        text_box(s, x=x, y=bar_y + 0.08, w=cw, h=0.42, text=num,
                 size=22, bold=True, color=nc, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
        text_box(s, x=x + 0.03, y=bar_y + 0.47, w=cw - 0.06, h=ch - 0.50,
                 text=sec_title, size=8.5, bold=(i == here_idx), color=tc,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.12)
    speaker_notes(s, load_notes(notes_slide_id))


# ============================================================
# Reusable content blocks
# ============================================================

def two_col_compare(s, *, box_y, box_h, left, right,
                    left_color=MID, right_color=TEAL, body_size=17):
    """left/right = dict(title, icon, icon_color, lines:[str]).
    Fixed even row slots — each line gets an equal vertical slot computed
    from box_h so text never overflows the ocean box. body_size raised to
    fill the taller box (P1-1 vertical-fill: ~85-90% canvas, projector-safe
    ≥16pt body, lec-07 baseline)."""
    box_x, box_w = 0.55, 12.25
    ocean_box(s, box_x, box_y, box_w, box_h)
    mid_x = box_x + box_w / 2
    filled_rect(s, mid_x - 0.01, box_y + 0.40, 0.02, box_h - 0.80,
                COVER_OUTLINE)
    header_h = 1.25
    for cx, cfg, ccol in [
        (box_x + 0.50, left, left_color),
        (mid_x + 0.50, right, right_color),
    ]:
        cw = box_w / 2 - 1.00
        add_image(s, icon(cfg["icon"], cfg.get("icon_color", "blue")),
                  cx, box_y + 0.40, 0.70, 0.70)
        text_box(s, x=cx + 0.92, y=box_y + 0.42, w=cw - 0.92, h=0.78,
                 text=cfg["title"], size=20, bold=True, color=ccol,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        lines = cfg["lines"]
        area_y = box_y + header_h
        area_h = box_h - header_h - 0.28
        slot = area_h / len(lines)
        for i, ln in enumerate(lines):
            sy = area_y + i * slot
            filled_rect(s, cx + 0.02, sy + slot / 2 - 0.06, 0.12, 0.12, ccol,
                        radius=True, radius_adj=0.5)
            text_box(s, x=cx + 0.32, y=sy, w=cw - 0.24, h=slot,
                     text=ln, size=body_size, color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.22)
    return box_x, box_w


def pipeline_row(s, *, y, h, stages, gold_idx=None):
    """stages = list of (title, body). gold_idx → gold-framed stage.
    Fonts raised for projector readability (P1-1 / Schema §5.5)."""
    total_w = 12.25
    x0 = 0.55
    n = len(stages)
    aw = 0.52
    sw = (total_w - (n - 1) * aw) / n
    head_h = 1.00
    for i, (title, body) in enumerate(stages):
        x = x0 + i * (sw + aw)
        is_gold = (i == gold_idx)
        ocean_box(s, x, y, sw, h,
                  fill=GOLD_TINT if is_gold else SURFACE,
                  stroke=GOLD if is_gold else LIGHT,
                  stroke_pt=2.0 if is_gold else 1.5)
        filled_rect(s, x + 0.18, y + 0.18, sw - 0.36, head_h,
                    GOLD if is_gold else MID, radius=True, radius_adj=0.14)
        text_box(s, x=x + 0.22, y=y + 0.18, w=sw - 0.44, h=head_h,
                 text=f"{i + 1}. {title}", size=15, bold=True,
                 color=DEEP if is_gold else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.10)
        text_box(s, x=x + 0.24, y=y + head_h + 0.30, w=sw - 0.48,
                 h=h - head_h - 0.48,
                 text=body, size=14, color=DEEP, line_spacing=1.28,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            right_arrow(s, x + sw + 0.05, y + h / 2 - 0.24, aw - 0.10, 0.48,
                        fill=MID)


# ============================================================
# Slide builders — 32 slides
# ============================================================

def build_s01(p):
    """Cover."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=7.6, y=1.7, w=5.7, h=5.5, text="06",
             size=320, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.7, y=1.05, w=7.0, h=0.55, text="ЛЕКЦИЯ",
             size=18, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    filled_rect(s, 0.72, 1.60, 0.7, 0.05, fill=TEAL)
    text_box(s, x=0.7, y=2.05, w=8.2, h=3.0,
             text="AI в инженерном\nпроектировании\nи CAD/CAM",
             size=52, bold=True, color=DEEP, line_spacing=1.06,
             align=PP_ALIGN.LEFT)
    filled_rect(s, 0.7, 5.55, 0.05, 0.55, fill=TEAL)
    text_box(s, x=0.95, y=5.55, w=10.8, h=0.6,
             text="Назови вид ИИ, прежде чем доверить ему проектное решение.",
             size=20, color=MID, align=PP_ALIGN.LEFT, line_spacing=1.25)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """Hook — кронштейн «какой вид ИИ»."""
    s = blank(p)
    slide_title(s, "Какой вид ИИ спроектировал этот кронштейн?", size=26)
    # Left: bionic/topology-optimized bracket in ocean box (~52%).
    # P1-4: recognizable organic "grown" lattice silhouette (2 mount hubs +
    # load eye + branching ribs + porous lightening holes) — drives the
    # cognitive dissonance "this looks AI-grown / no, it is not AI".
    bx, by, bw, bh = 0.55, 1.55, 6.6, 4.85
    ocean_box(s, bx, by, bw, bh)
    add_image(s, ASSETS / "illustrations/bionic-bracket.png",
              bx + 0.30, by + 0.22, bw - 0.6, bh - 0.92)
    text_box(s, x=bx, y=by + bh - 0.52, w=bw, h=0.42,
             text="«Органическая» ажурная форма — выглядит как «выращенная» "
                  "нейросетью",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Right: reveal block
    rx, ry, rw, rh = 7.35, 1.55, 5.45, 4.85
    ocean_box(s, rx, ry, rw, rh)
    text_box(s, x=rx + 0.40, y=ry + 0.45, w=rw - 0.8, h=0.45,
             text="Ответ:", size=16, bold=True, color=TEAL)
    text_runs(s, rx + 0.40, ry + 1.05, rw - 0.8, 3.10, [
        {"text": "Ни один\nгенеративный AI.", "size": 28, "color": DEEP,
         "bold": True},
        {"newpara": True, "text": "Детерминированный\nчисленный оптимизатор.",
         "size": 21, "color": MID, "bold": True, "line_spacing": 1.22},
        {"newpara": True, "text": "Постановка — ", "size": 21,
         "color": DEEP, "bold": True, "line_spacing": 1.22},
        {"text": "1904 год", "size": 25, "color": GOLD, "bold": True},
        {"text": ".", "size": 21, "color": DEEP, "bold": True},
    ], line_spacing=1.32)
    gold_callout(s, rx + 0.40, ry + rh - 1.05, rw - 0.8, 0.80,
                 "8 деталей → 1  ·  масса −~40%\n>150 вариантов перебора",
                 size=14, align=PP_ALIGN.CENTER)
    footer(s, "GM seat bracket (General Motors × Autodesk, 2018) · постановка "
              "восходит к работе Мичелла, 1904.")
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    """Lecture map — 3 вопроса + 7-card roadmap."""
    s = blank(p)
    slide_title(s, "Спайн лекции — три вопроса к самому себе и шесть частей.",
                size=25)
    qx, qy, qw, qh = 0.55, 1.45, 12.25, 2.95
    ocean_box(s, qx, qy, qw, qh)
    text_box(s, x=qx + 0.35, y=qy + 0.22, w=qw - 0.7, h=0.40,
             text="Три вопроса-якоря — повторяются в каждой части",
             size=14, bold=True, color=MID)
    qs = [
        ("1", "Детерминированная математика или обученная вероятностная "
              "модель?"),
        ("2", "Что оптимизируется и кто задал ограничения?"),
        ("3", "Кто несёт ответственность?"),
    ]
    iy = qy + 0.78
    for num, q in qs:
        filled_rect(s, qx + 0.35, iy, 0.48, 0.48, MID, radius=True,
                    radius_adj=0.5)
        text_box(s, x=qx + 0.35, y=iy, w=0.48, h=0.48, text=num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=qx + 0.98, y=iy + 0.03, w=qw - 1.4, h=0.55, text=q,
                 size=18, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.12)
        iy += 0.66
    # 7-card roadmap (part 0 active)
    bar_y = 4.70
    n = 7
    gap = 0.12
    cw = (12.3 - gap * (n - 1)) / n
    sx = 0.55
    ch = 1.55
    parts = [
        ("0", "Старт\n+ карта"),
        ("1", "Фундамент\n+ таксономия"),
        ("2", "Генеративный\nдизайн ≠ genAI"),
        ("3", "Суррогат\n/ PINN"),
        ("4", "Генеративный\nAI / LLM"),
        ("5", "Когда\nсказать «нет»"),
        ("6", "Синтез"),
    ]
    for i, (num, t) in enumerate(parts):
        x = sx + i * (cw + gap)
        active = (i == 0)
        ocean_box(s, x, bar_y, cw, ch,
                  fill=GOLD if active else WHITE,
                  stroke=GOLD if active else LIGHT,
                  stroke_pt=2.0 if active else 1.0)
        text_box(s, x=x, y=bar_y + 0.18, w=cw, h=0.55, text=num,
                 size=26, bold=True, color=WHITE if active else MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        text_box(s, x=x + 0.04, y=bar_y + 0.78, w=cw - 0.08, h=ch - 0.85,
                 text=t, size=10, bold=active,
                 color=WHITE if active else SLATE,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)
    footer(s, "Классификация — для всех 6 классов ИИ; глубина — для трёх с "
              "высшей ценой ошибки.")
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    build_section_divider(
        p, here_idx=1, title="Назови\nвид ИИ",
        frame_phrase="фундамент — решать уравнение vs угадывать по примерам — "
                     "и таксономия шести классов",
        notes_slide_id="s04")


def build_s05(p):
    """Фундамент: решает уравнение vs угадывает."""
    s = blank(p)
    slide_title(s, "Детерминированная оптимизация решает уравнение; обученная "
                   "модель угадывает по примерам.", size=24, h=1.15)
    left = {
        "title": "РЕШАЕТ заданное уравнение",
        "icon": "function-square", "icon_color": "blue",
        "lines": [
            "Инженер формулирует цель и ограничения (равновесие, прочность, "
            "объём)",
            "Алгоритм численно ищет точку, удовлетворяющую уравнениям",
            "Физический закон присутствует явно",
            "Тот же вход → тот же ответ · воспроизводимо и прослеживаемо",
        ],
    }
    right = {
        "title": "УГАДЫВАЕТ по примерам",
        "icon": "dices", "icon_color": "teal",
        "lines": [
            "Модель обучена на корпусе примеров",
            "Воспроизводит «как обычно выглядит / звучит ответ»",
            "Уравнения равновесия внутри нет",
            "Ошибка не самодиагностируется",
        ],
    }
    two_col_compare(s, box_y=1.55, box_h=4.55, left=left, right=right)
    gold_callout(s, 0.55, 6.25, 12.25, 0.78,
                 "Любой инструмент «с ИИ» сначала располагают на этой оси.",
                 size=19, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """Таксономия — одна ось, 6 ярлыков (schema_axis)."""
    s = blank(p)
    slide_title(s, "Шесть классов ИИ ложатся на одну ось: чем правее — тем "
                   "меньше встроенных гарантий.", size=24, h=1.05)
    # Visible "reference skeleton — don't memorize now" signal (P1-6).
    sig_x, sig_y, sig_w, sig_h = 0.55, 1.42, 12.25, 0.56
    filled_rect(s, sig_x, sig_y, sig_w, sig_h, GOLD_TINT, stroke=GOLD,
                stroke_pt=1.4, radius=True, radius_adj=0.18)
    add_image(s, icon("bookmark", "gold"), sig_x + 0.22, sig_y + 0.10,
              0.36, 0.36)
    text_box(s, x=sig_x + 0.72, y=sig_y, w=sig_w - 0.95, h=sig_h,
             text="Это справочный скелет — шесть классов не нужно заучивать "
                  "сразу: они наполнятся смыслом по ходу лекции.",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.10)
    box_x, box_y, box_w, box_h = 0.55, 2.14, 12.25, 4.30
    ocean_box(s, box_x, box_y, box_w, box_h)
    # Axis line with arrow
    axis_y = box_y + 0.90
    filled_rect(s, box_x + 0.55, axis_y, box_w - 2.0, 0.06, MID)
    right_arrow(s, box_x + box_w - 1.45, axis_y - 0.13, 0.55, 0.32, fill=GOLD)
    text_box(s, x=box_x + 0.45, y=box_y + 0.28, w=4.0, h=0.40,
             text="детерминированное", size=14, bold=True, color=MID)
    text_box(s, x=box_x + box_w - 5.0, y=box_y + 0.28, w=4.5, h=0.40,
             text="вероятностное →  меньше гарантий", size=14, bold=True,
             color=GOLD, align=PP_ALIGN.RIGHT)
    classes = [
        ("sliders", "Оптим. ML /\nтопологическая\nоптимизация",
         "двигает плотность материала,\nмин. массу при σ ≤ [σ]", MID),
        ("git-branch", "Эволюционные\n/ GA",
         "мутация и отбор по\nзаданной целевой функции", MID),
        ("gauge", "Суррогат\n/ PINN",
         "нейросеть выучила вход→\nвыход тысяч расчётов", LIGHT),
        ("scan-eye", "Computer\nVision",
         "вероятностный\nклассификатор изображений", LIGHT),
        ("message-square", "LLM-\nассистент",
         "текстовые черновики,\nобъяснение пунктов норм", TEAL),
        ("box", "Генеративный AI\nгеометрии",
         "порождает 3D без\nфизических гарантий", TEAL),
    ]
    n = 6
    cw = (box_w - 0.9) / n
    cx0 = box_x + 0.45
    card_y = axis_y + 0.42
    for i, (ic, name, defn, col) in enumerate(classes):
        x = cx0 + i * cw
        # marker dot on axis
        filled_rect(s, x + cw / 2 - 0.08, axis_y - 0.06, 0.18, 0.18, col,
                    radius=True, radius_adj=0.5)
        add_image(s, icon(ic, "blue"), x + cw / 2 - 0.28, card_y, 0.56, 0.56)
        text_box(s, x=x + 0.03, y=card_y + 0.64, w=cw - 0.06, h=0.95,
                 text=f"{i + 1}. {name}", size=13, bold=True, color=col,
                 align=PP_ALIGN.CENTER, line_spacing=1.08)
        text_box(s, x=x + 0.03, y=card_y + 1.62, w=cw - 0.06, h=1.05,
                 text=defn, size=12, color=DEEP, align=PP_ALIGN.CENTER,
                 line_spacing=1.16)
    text_box(s, x=box_x + 0.45, y=box_y + box_h - 0.46, w=box_w - 0.9, h=0.38,
             text="классы 1–2 — оптимизация заданной функции · 3 — "
                  "аппроксимация решателя · 4–6 — вероятностные модели",
             size=12.5, italic=True, color=SLATE, align=PP_ALIGN.CENTER)
    footer(s, "Полная атрибутивная матрица — в финале лекции, когда каждая "
              "ячейка наполнится смыслом.")
    speaker_notes(s, load_notes("s06"))


def _matrix_da_net(s, *, box_y, box_h, rows, col_titles, gold_col=1,
                   body_size=12, label_size=13):
    """Generic ДА/НЕТ/альтернатива matrix (canonical format-template grid).
    rows = list of (icon, row_title, [c1, c2, c3], row_color)."""
    bx, bw = 0.55, 12.25
    ocean_box(s, bx, box_y, bw, box_h)
    label_w = 2.75
    col_w = (bw - label_w - 0.6) / 3
    # header row
    hy = box_y + 0.22
    for c, ct in enumerate(col_titles):
        cx = bx + label_w + 0.30 + c * col_w
        fill = GOLD if c == gold_col else MID
        filled_rect(s, cx, hy, col_w - 0.12, 0.52, fill, radius=True,
                    radius_adj=0.16)
        text_box(s, x=cx + 0.05, y=hy, w=col_w - 0.22, h=0.52, text=ct,
                 size=14, bold=True, color=DEEP if c == gold_col else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    row_y = hy + 0.68
    rh = (box_h - 0.90 - 0.24) / len(rows)
    for ri, (ic, rt, cells, rcol) in enumerate(rows):
        ry = row_y + ri * rh
        # row label cell
        filled_rect(s, bx + 0.25, ry + 0.07, label_w, rh - 0.14,
                    TEAL_TINT if rcol == TEAL else SURFACE,
                    stroke=rcol, stroke_pt=1.2, radius=True, radius_adj=0.10)
        add_image(s, icon(ic, "blue"), bx + 0.42, ry + rh / 2 - 0.26,
                  0.50, 0.50)
        text_box(s, x=bx + 1.00, y=ry + 0.07, w=label_w - 0.78, h=rh - 0.14,
                 text=rt, size=label_size, bold=True, color=rcol,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
        for c, txt in enumerate(cells):
            cx = bx + label_w + 0.30 + c * col_w
            is_gold = (c == gold_col)
            filled_rect(s, cx, ry + 0.07, col_w - 0.12, rh - 0.14,
                        GOLD_TINT if is_gold else WHITE,
                        stroke=GOLD if is_gold else LIGHT,
                        stroke_pt=1.2 if is_gold else 1.0,
                        radius=True, radius_adj=0.10)
            text_box(s, x=cx + 0.16, y=ry + 0.07, w=col_w - 0.42,
                     h=rh - 0.14, text=txt, size=body_size, color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
    return bx, bw


def build_s07(p):
    """Пара 1 — матрица 2×3 (канонический формат-шаблон: грид ДА/НЕТ/альт).
    Дифференцирована от s08 (s08 = вертикальные класс-карточки)."""
    s = blank(p)
    slide_title(s, "Оптимизационный ML и суррогат/PINN — узнать по границе.",
                size=25)
    _matrix_da_net(
        s, box_y=1.55, box_h=4.85,
        col_titles=["Когда ДА", "Когда НЕТ", "Классическая альтернатива"],
        rows=[
            ("sliders", "Оптимизационный ML / топологическая оптимизация",
             ["снижение массы при заданной прочности; детали под аддитив",
              "ограничения заданы неполно — оптимум под неверную постановку",
              "параметрическая оптимизация + расчёт по нормам"], MID),
            ("gauge", "Суррогат / PINN",
             ["быстрый прогон многих вариантов внутри обучающей области",
              "экстраполяция за обучение, разрывы полей, сертификация",
              "полный МКЭ / CFD с проверкой сходимости сетки"], LIGHT),
        ], gold_col=1, body_size=14, label_size=15)
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """Пара 2 — ТРИ вертикальные класс-карточки (умышленно ИНОЙ layout, чем
    s07-грид: разбивает монотонность s06→s07→s08). Каждая карточка:
    цветной хедер с иконкой + 3 сегмента ✓ДА / ✗НЕТ / ↳альтернатива."""
    s = blank(p)
    slide_title(s, "Генеративный AI/LLM, Computer Vision и эволюционные/GA — "
                   "по одному примеру и границе.", size=23, h=1.05)
    cards = [
        ("message-square", "Генеративный AI / LLM-ассистент", TEAL,
         "концепт-скетч, мудборд; черновик текста под проверку",
         "деталь в расчёт/производство без переработки; истина по нормам",
         "параметрика + детерм. топ-оптимизация + КЭ; норматив"),
        ("scan-eye", "Computer Vision", LIGHT,
         "высокообъёмная рутинная сортировка как первый фильтр",
         "единственный арбитр годности в безопасностно-критичном НК",
         "нормированный НК с дефектоскопистом и методикой POD "
         "(вероятность обнаружения дефекта)"),
        ("git-branch", "Эволюционные / GA", MID,
         "противоречивые требования (антенна NASA ST5)",
         "нужна сертификационная объяснимость и интерпретируемость",
         "аналитический / градиентный синтез по методике"),
    ]
    cw = 4.05
    gap = 0.15
    cy, ch = 1.55, 4.85
    seg_defs = [("✓  когда ДА", GOLD, GOLD_TINT),
                ("✗  когда НЕТ", LIGHT, WHITE),
                ("↳  альтернатива", MID, TEAL_TINT)]
    for i, (ic, name, accent, da, net, alt) in enumerate(cards):
        x = 0.55 + i * (cw + gap)
        ocean_box(s, x, cy, cw, ch, stroke=accent, stroke_pt=1.8)
        # colored header band — distinct visual signature vs s07 grid
        filled_rect(s, x + 0.18, cy + 0.18, cw - 0.36, 0.98, accent,
                    radius=True, radius_adj=0.14)
        add_image(s, icon(ic, "white"), x + 0.34, cy + 0.38, 0.60, 0.60)
        text_box(s, x=x + 1.04, y=cy + 0.20, w=cw - 1.22, h=0.94, text=name,
                 size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.10)
        seg_y = cy + 1.32
        seg_h = (ch - 1.32 - 0.20) / 3
        for (lab, sc, sf), body in zip(seg_defs, (da, net, alt)):
            filled_rect(s, x + 0.18, seg_y, cw - 0.36, seg_h - 0.12, sf,
                        stroke=sc, stroke_pt=1.0, radius=True,
                        radius_adj=0.07)
            text_box(s, x=x + 0.34, y=seg_y + 0.09, w=cw - 0.62, h=0.30,
                     text=lab, size=12.5, bold=True,
                     color=DEEP if sc == GOLD else sc)
            text_box(s, x=x + 0.34, y=seg_y + 0.40, w=cw - 0.62,
                     h=seg_h - 0.54, text=body, size=12, color=DEEP,
                     anchor=MSO_ANCHOR.TOP, line_spacing=1.16)
            seg_y += seg_h
    gold_callout(s, 0.55, 6.55, 12.25, 0.62,
                 "Якоря: детерм./вероятн.?  ·  что оптимизируется, кто задал "
                 "ограничения?  ·  кто отвечает?", size=14,
                 align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    build_section_divider(
        p, here_idx=2, title="Генеративный дизайн —\nэто не генеративный AI",
        frame_phrase="детерминированный оптимизатор на классической "
                     "математике (Мичелл, Коши, SIMP), а не порождающая "
                     "модель",
        notes_slide_id="s09")


def build_s10(p):
    """Терминологическая ловушка — 2 колонки контраст."""
    s = blank(p)
    slide_title(s, "«AI generates the design» — маркетинговая подмена понятий.",
                size=25)
    left = {
        "title": "Маркетинг говорит",
        "icon": "megaphone", "icon_color": "teal",
        "lines": [
            "«AI generates the design»",
            "«ИИ придумал деталь»",
            "«творческий искусственный интеллект»",
        ],
    }
    right = {
        "title": "Что на самом деле",
        "icon": "cog", "icon_color": "blue",
        "lines": [
            "Топологическая оптимизация — детерминированный численный метод",
            "Тот же вход → тот же ответ",
            "Ни одной обученной порождающей модели",
            "Сам вендор различает топ-оптимизацию и «генеративный AI»",
        ],
    }
    two_col_compare(s, box_y=1.55, box_h=4.55, left=left, right=right)
    gold_callout(s, 0.55, 6.25, 12.25, 0.78,
                 "Урок: прежде чем доверять — спроси, какой это вид ИИ.",
                 size=19, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s10"))


def build_s11(p):
    """Классическая родословная → SIMP (schema_timeline)."""
    s = blank(p)
    slide_title(s, "Топологическая оптимизация — это градиентный спуск по "
                   "плотности материала.", size=24, h=1.15)
    # Visible "lineage, not chronology" annotation (P2: 1904→1847→1988→1989
    # is a logical родословная, not a date sequence — pre-empt misread).
    ann_x, ann_y, ann_w, ann_h = 0.55, 1.42, 12.25, 0.54
    filled_rect(s, ann_x, ann_y, ann_w, ann_h, TEAL_TINT, stroke=TEAL,
                stroke_pt=1.2, radius=True, radius_adj=0.18)
    add_image(s, icon("git-branch", "teal"), ann_x + 0.22, ann_y + 0.10,
              0.34, 0.34)
    text_box(s, x=ann_x + 0.70, y=ann_y, w=ann_w - 0.95, h=ann_h,
             text="Это родословная (логический порядок), а не хронология: "
                  "постановка Мичелла (1904) — основа; алгоритм Коши (1847) "
                  "старше по дате, но применяется к ней позже.",
             size=13, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.08)
    box_x, box_y, box_w, box_h = 0.55, 2.10, 12.25, 4.32
    ocean_box(s, box_x, box_y, box_w, box_h)
    nodes = [
        ("Известно из\nсопромата", "недогруженный материал\nможно убрать "
         "(равнопрочность)", MID, False),
        ("Мичелл 1904", "интуиция и постановка\n(минимум материала)", MID,
         False),
        ("Коши 1847 /\nградиентный спуск", "тот же алгоритм\nнаискорейшего "
         "спуска", LIGHT, False),
        ("Бендсё–Кикучи\n1988", "численный метод:\nнепрерывная плотность",
         LIGHT, False),
        ("SIMP 1989", "индустриальный стандарт:\nρ∈[0,1], штраф к «ч/б»",
         GOLD, True),
    ]
    n = len(nodes)
    seg = (box_w - 1.1) / n
    # baseline through vertical centre of box (label above, sub below)
    base_y = box_y + 1.95
    filled_rect(s, box_x + 0.55, base_y, box_w - 1.1, 0.05, LIGHT)
    text_box(s, x=box_x + 0.55, y=box_y + 0.12, w=box_w - 1.1, h=0.32,
             text="логическая родословная  →", size=12, bold=True,
             italic=True, color=LIGHT, align=PP_ALIGN.RIGHT)
    for i, (title, sub, col, is_pivot) in enumerate(nodes):
        cx = box_x + 0.55 + seg * i + seg / 2
        # connector arrow segment (behind dot)
        if i < n - 1:
            right_arrow(s, cx + 0.18, base_y - 0.085, seg - 0.36, 0.22,
                        fill=COVER_OUTLINE)
        dr = 0.32 if is_pivot else 0.18
        filled_rect(s, cx - dr / 2, base_y + 0.025 - dr / 2, dr, dr, col,
                    radius=True, radius_adj=0.5)
        # title ABOVE node, close to it
        text_box(s, x=cx - seg / 2 + 0.06, y=box_y + 0.50,
                 w=seg - 0.12, h=1.20,
                 text=title, size=16 if is_pivot else 14,
                 bold=True, color=GOLD if is_pivot else DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
                 line_spacing=1.10)
        # sub BELOW node, close to it
        text_box(s, x=cx - seg / 2 + 0.06, y=base_y + 0.32,
                 w=seg - 0.12, h=1.15,
                 text=sub, size=13, color=SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                 line_spacing=1.16)
    gold_callout(s, box_x + 0.5, box_y + box_h - 0.92, box_w - 1.0, 0.70,
                 "на каждом шаге явно решается уравнение равновесия — "
                 "сходимость к оптимуму, НЕ сэмпл из распределения",
                 size=14, align=PP_ALIGN.CENTER)
    footer(s, "KKT / множители Лагранжа — формальное условие оптимума; "
              "формул на слайде нет.")
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """Реальные кейсы GM + Airbus — 2 карточки (cascade ≤3 числа)."""
    s = blank(p)
    slide_title(s, "GM и Airbus: радикальное снижение массы — численная "
                   "оптимизация и перебор, не diffusion.", size=23, h=1.15)
    cards = [
        ("car", "GM seat bracket",
         ["8 деталей → 1", "масса −~40%", ">150 вариантов перебора"],
         "вид ИИ: численная оптимизация + облачный перебор"),
        ("plane", "Airbus bionic partition",
         ["масса −~45%", "≈35 кг против ≈65 кг", "материал Scalmalloy"],
         "вид ИИ: топологическая/решётчатая оптимизация + перебор"),
    ]
    cw = 6.0
    gap = 0.25
    cx0 = 0.55
    cy, ch = 1.60, 4.55
    for i, (ic, title, nums, label) in enumerate(cards):
        x = cx0 + i * (cw + gap)
        ocean_box(s, x, cy, cw, ch)
        add_image(s, icon(ic, "blue"), x + 0.40, cy + 0.40, 0.95, 0.95)
        text_box(s, x=x + 1.50, y=cy + 0.52, w=cw - 1.70, h=0.75,
                 text=title, size=22, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        ny = cy + 1.75
        for num in nums:
            filled_rect(s, x + 0.45, ny + 0.10, 0.14, 0.14, MID,
                        radius=True, radius_adj=0.5)
            text_box(s, x=x + 0.78, y=ny, w=cw - 1.1, h=0.50, text=num,
                     size=19, bold=True, color=DEEP, line_spacing=1.10)
            ny += 0.66
        filled_rect(s, x + 0.45, cy + ch - 1.05, cw - 0.90, 0.78, MID,
                    radius=True, radius_adj=0.14)
        text_box(s, x=x + 0.62, y=cy + ch - 1.05, w=cw - 1.24, h=0.78,
                 text=label, size=14, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.16)
    gold_callout(s, 0.55, 6.32, 12.25, 0.72,
                 "Обе детали требуют аддитива — это ограничение, не "
                 "случайность.", size=16, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """Россия honest — 3 карточки."""
    s = blank(p)
    slide_title(s, "В России APM FEM в КОМПАС-3D — это детерминированный SIMP, "
                   "а не AI.", size=24, h=1.05)
    cards = [
        ("award", "Реальные достижения", MID, [
            "ОКБ Сухого — силовой кронштейн Су-57 ~на четверть легче",
            "цифровая модель на суперкомпьютере · LPBF "
            "(лазерное послойное сплавление)",
            "прорабатывался прототип для МС-21-300",
        ]),
        ("alert-triangle", "Подмена понятий", GOLD, [
            "APM FEM исключает существенную долю объёма с сохранением "
            "прочности",
            "но это детерминированный SIMP, НЕ AI — та же ловушка",
            "ядро C3D — стратегический актив (в мире — единицы ядер)",
        ]),
        ("circle-slash", "Где честно отстаём", LIGHT, [
            "зрелого AI-генеративного движка уровня Fusion в РФ нет",
            "честная альтернатива — классическая топ-оптимизация",
            "+ ручная реконструкция в редактируемое тело",
        ]),
    ]
    cw = 4.05
    gap = 0.15
    cy, ch = 1.45, 5.30
    for i, (ic, title, tc, lines) in enumerate(cards):
        x = 0.55 + i * (cw + gap)
        is_gold = (tc == GOLD)
        ocean_box(s, x, cy, cw, ch,
                  fill=GOLD_TINT if is_gold else SURFACE,
                  stroke=GOLD if is_gold else LIGHT,
                  stroke_pt=2.0 if is_gold else 1.5)
        add_image(s, icon(ic, "gold" if is_gold else "blue"),
                  x + 0.32, cy + 0.34, 0.88, 0.88)
        text_box(s, x=x + 0.32, y=cy + 1.38, w=cw - 0.64, h=0.85,
                 text=title, size=18, bold=True,
                 color=GOLD if is_gold else tc, line_spacing=1.12)
        ly = cy + 2.35
        for ln in lines:
            filled_rect(s, x + 0.32, ly + 0.09, 0.11, 0.11,
                        GOLD if is_gold else MID, radius=True, radius_adj=0.5)
            text_box(s, x=x + 0.54, y=ly, w=cw - 0.80, h=0.92, text=ln,
                     size=13.5, color=DEEP, line_spacing=1.24)
            ly += 1.02
    footer(s, "Точная доля объёма, исключаемого APM FEM, — по заявлению "
              "вендора.")
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """Провал раздела — pipeline garbage-in → optimal garbage."""
    s = blank(p)
    slide_title(s, "Главный риск топологической оптимизации — не "
                   "«галлюцинация», а garbage-in → optimal garbage.",
                   size=24, h=1.15)
    pipeline_row(
        s, y=1.60, h=3.05,
        stages=[
            ("Неполная постановка",
             "забыт нагрузочный случай;\nне заданы технологические\n"
             "ограничения"),
            ("Детерминированный\nоптимизатор",
             "честно решает\nпоставленную задачу"),
            ("Безупречный оптимум —\nинженерно непригодная форма",
             "математически верна,\nизготовить нельзя"),
        ], gold_idx=2)
    # callout + alternative
    cx, cy, cwd, chd = 0.55, 4.90, 6.05, 2.10
    ocean_box(s, cx, cy, cwd, chd)
    add_image(s, icon("ban", "blue"), cx + 0.30, cy + 0.34, 0.84, 0.84)
    text_box(s, x=cx + 1.28, y=cy + 0.26, w=cwd - 1.54, h=chd - 0.52,
             text="Неизготовимость: замкнутые полости и криволинейные рёбра "
                  "без доступа фрезы — только аддитив; трудно для обмера и "
                  "нормоконтроля по ЕСКД.",
             size=14, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.24)
    gold_callout(s, 6.75, 4.90, 6.05, 2.10,
                 "Альтернатива для серийной фрезеровки/литья — классическая "
                 "параметрика + ручная топологическая оптимизация под "
                 "технологичность.",
                 size=15)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    build_section_divider(
        p, here_idx=3, title="ИИ ускоряет счёт,\nне заменяет физику",
        frame_phrase="суррогат угадывает выход решателя по примерам; "
                     "PINN сглаживает скачки",
        notes_slide_id="s15")


def build_s16(p):
    """КЭ vs суррогат — 2 колонки."""
    s = blank(p)
    slide_title(s, "КЭ-анализ решает уравнения равновесия; суррогат угадывает "
                   "поле по примерам.", size=24, h=1.15)
    left = {
        "title": "КЭ-решатель",
        "icon": "grid-3x3", "icon_color": "blue",
        "lines": [
            "деталь → сетка → решает систему уравнений равновесия → поле "
            "напряжений",
            "проверенная, прослеживаемая, детерминированная математика",
            "контроль погрешности проверкой сходимости сетки",
        ],
    }
    right = {
        "title": "Суррогат",
        "icon": "gauge", "icon_color": "teal",
        "lines": [
            "выучил отображение «геометрия+нагрузки → напряжения» на тысячах "
            "расчётов",
            "НЕ решает уравнение — угадывает по примерам",
            "обучение оффлайн, заранее, на дорогом архиве",
        ],
    }
    two_col_compare(s, box_y=1.55, box_h=4.55, left=left, right=right)
    gold_callout(s, 0.55, 6.25, 12.25, 0.78,
                 "Под капотом — система уравнений равновесия; суррогат её не "
                 "решает.", size=18, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """Состояние 2026 — 3 продукт-карточки (cascade ≤4 названия, диапазоны)."""
    s = blank(p)
    slide_title(s, "Суррогаты-2026: «до N×» — верхняя граница, не типичное.",
                size=25)
    subhead(s, "К 2026 крупные вендоры превратили суррогаты в продукт — "
               "числа читаем критически.")
    cards = [
        ("Ansys SimAI", "релиз 2026 R1",
         "суррогат на основе редуцированной модели; ускорение — от десятков "
         "до сотен раз"),
        ("Altair PhysicsAI", "обучение на архиве",
         "учится на историческом архиве симуляций; ускорение до порядка "
         "тысячи раз"),
        ("NVIDIA PhysicsNeMo", "open-source",
         "открытый Python-фреймворк physics-AI: от PINN до нейронных "
         "операторов"),
    ]
    cw = 4.05
    gap = 0.15
    cy, ch = 1.80, 3.50
    for i, (name, tag, body) in enumerate(cards):
        x = 0.55 + i * (cw + gap)
        ocean_box(s, x, cy, cw, ch)
        add_image(s, icon("zap", "blue"), x + 0.32, cy + 0.34, 0.74, 0.74)
        text_box(s, x=x + 1.20, y=cy + 0.36, w=cw - 1.40, h=0.68,
                 text=name, size=18, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        text_box(s, x=x + 0.32, y=cy + 1.24, w=cw - 0.64, h=0.38, text=tag,
                 size=13, italic=True, color=TEAL)
        text_box(s, x=x + 0.32, y=cy + 1.70, w=cw - 0.64, h=1.65, text=body,
                 size=14.5, color=DEEP, line_spacing=1.30)
    gold_callout(s, 0.55, 5.55, 12.25, 1.48,
                 "«до ~10²–10³×» (≈ ×100–×1000) — на выгодном для вендора "
                 "частном случае. Стоимость генерации обучающего архива "
                 "(тысячи дорогих прогонов решателя) из ROI исключают.",
                 size=16)
    speaker_notes(s, load_notes("s17"))


def build_s18(p):
    """PINN размазывает концентратор — эскиз-график + карточка."""
    s = blank(p)
    slide_title(s, "PINN «размазывает» пик у концентратора и выдаёт заниженное "
                   "напряжение.", size=24, h=1.15)
    # Left: stress-curve sketch — rebuilt for projector legibility (P1-7).
    # Two SOLID filled freeform curves (no thin staircase polylines), labels
    # placed in clear space away from the curves, single bold gold gap
    # annotation between the peak tops.
    gx, gy, gw, gh = 0.55, 1.55, 6.85, 4.85
    ocean_box(s, gx, gy, gw, gh)
    import math
    # plot frame
    ox = gx + 0.55              # y-axis x
    oy = gy + gh - 1.05         # x-axis y (baseline)
    plot_top = gy + 1.05
    span_w = gw - 1.20
    filled_rect(s, ox, plot_top, 0.035, oy - plot_top, SLATE)  # y-axis
    filled_rect(s, ox, oy, span_w + 0.20, 0.035, SLATE)        # x-axis
    text_box(s, x=gx + 0.30, y=gy + 0.22, w=gw - 0.6, h=0.34,
             text="напряжение вдоль сечения у отверстия / выточки",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

    base_lvl = 0.55             # nominal far-field stress height (in)
    peak_h = 2.45               # true sharp peak height
    pinn_h = 1.30               # PINN smoothed (lower) peak height

    def curve_pts(amp, sigma):
        pts = []
        for k in range(81):
            t = k / 80.0
            xx = ox + 0.12 + t * span_w
            yy = oy - (base_lvl + amp * math.exp(-((t - 0.5) ** 2) / sigma))
            pts.append((xx, yy))
        return pts

    def filled_curve(pts, color, line_color):
        fb = s.shapes.build_freeform(Inches(pts[0][0]), Inches(oy))
        fb.add_line_segments(
            [(Inches(x), Inches(y)) for x, y in pts] +
            [(Inches(pts[-1][0]), Inches(oy))], close=True)
        shp = fb.convert_to_shape()
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = line_color
        shp.line.width = Pt(2.0)
        disable_shadow(shp)
        return shp

    # PINN (smoothed) drawn first, behind; teal tint area
    filled_curve(curve_pts(pinn_h, 0.020), TEAL_TINT, TEAL)
    # True sharp peak in front; light-blue area with deep outline
    filled_curve(curve_pts(peak_h, 0.0042), RGBColor(0xDD, 0xEC, 0xF2), MID)

    cx_peak = ox + 0.12 + 0.5 * span_w
    true_top = oy - (base_lvl + peak_h)
    pinn_top = oy - (base_lvl + pinn_h)
    # gold double-headed gap arrow between the two peak tops, placed in the
    # clear band well right of the spike (iter2: more clearance, no overlap).
    gap_x = cx_peak + 1.55
    up = s.shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW,
                            Inches(gap_x - 0.14), Inches(true_top + 0.04),
                            Inches(0.28), Inches(pinn_top - true_top - 0.08))
    up.fill.solid(); up.fill.fore_color.rgb = GOLD
    up.line.fill.background(); disable_shadow(up)
    text_box(s, x=gap_x + 0.26, y=(true_top + pinn_top) / 2 - 0.36,
             w=1.85, h=0.82,
             text="занижение\nнапряжения", size=13, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
    # small marker dot just below the spike apex (does not sit on the tip)
    filled_rect(s, cx_peak - 0.06, true_top + 0.16, 0.12, 0.12, MID,
                radius=True, radius_adj=0.5)
    # label: true peak (upper-left clear zone) + thin leader to the spike
    text_box(s, x=gx + 0.42, y=gy + 0.62, w=2.45, h=0.78,
             text="истинный пик —\nконцентратор", size=14, bold=True,
             color=MID, line_spacing=1.10)
    # label: PINN curve — sits over its own broad teal shoulder (right side,
    # mid-height) so the association reads without a leader line
    text_box(s, x=gx + gw - 2.75, y=oy - 1.55, w=2.45, h=0.78,
             text="PINN — размазанный,\nзаниженный пик", size=14, bold=True,
             color=TEAL, line_spacing=1.10, align=PP_ALIGN.RIGHT)
    # takeaway strip under x-axis
    text_box(s, x=gx + 0.40, y=oy + 0.16, w=gw - 0.80, h=0.74,
             text="заниженное напряжение в опасном сечении = "
                  "недооценённый риск разрушения", size=13, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER, line_spacing=1.16)
    # Right: other limitations + niche
    rx, ry, rw, rh = 7.60, 1.55, 5.20, 4.85
    ocean_box(s, rx, ry, rw, rh)
    text_box(s, x=rx + 0.32, y=ry + 0.28, w=rw - 0.64, h=0.48,
             text="Прочие ограничения", size=17, bold=True, color=DEEP)
    lims = [
        "не превосходит зрелые решатели на прямых задачах, часто медленнее",
        "не обобщается за пределы обучения",
        "при сбое непрозрачен",
        "усиливает высокочастотный шум данных",
    ]
    ly = ry + 1.00
    for ln in lims:
        filled_rect(s, rx + 0.32, ly + 0.08, 0.12, 0.12, LIGHT,
                    radius=True, radius_adj=0.5)
        text_box(s, x=rx + 0.56, y=ly, w=rw - 0.86, h=0.72, text=ln,
                 size=14, color=DEEP, line_spacing=1.22)
        ly += 0.82
    gold_callout(s, rx + 0.32, ry + rh - 1.30, rw - 0.64, 1.10,
                 "Где PINN реально полезен: разреженные данные, обратные "
                 "задачи, сложная геометрия — НЕ замена FEM/CFD.", size=14)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    """Зачем суррогат — pipeline + критерий (cascade ≤3 числа/4 названия)."""
    s = blank(p)
    slide_title(s, "Суррогат — для перебора; финал считает аттестованный "
                   "решатель.", size=25)
    pipeline_row(
        s, y=1.60, h=2.45,
        stages=[
            ("Сотни вариантов\nна этапе поиска",
             "а если ребро толще,\nа если другой материал"),
            ("Суррогат\nотсеивает",
             "прогон за минуты,\nотсев заведомо слабых"),
            ("5–10\nкандидатов", "узкий шорт-лист"),
            ("Аттестованный\nрешатель — финал",
             "детерминированный\nсертифицируемый расчёт"),
        ], gold_idx=3)
    gold_callout(s, 0.55, 4.30, 12.25, 1.05,
                 "Критерий: суррогат валиден только внутри обучающей области. "
                 "Экстраполяция = молчаливая ошибка (модель не сообщает, что "
                 "вышла за границу).", size=15)
    box_x, box_y, box_w, box_h = 0.55, 5.55, 12.25, 0.95
    ocean_box(s, box_x, box_y, box_w, box_h)
    text_box(s, x=box_x + 0.30, y=box_y + 0.10, w=box_w - 0.6, h=box_h - 0.20,
             text="Россия: «Логос» (Росатом) — суррогаты в обёртке, не в ядре "
                  "решателя · CML-Bench (СПбПУ) — накопленный архив → "
                  "ассистент.",
             size=12, italic=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.18)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    build_section_divider(
        p, here_idx=4, title="Где галлюцинация\nстоит дорого",
        frame_phrase="генеративный AI и LLM не содержат физического закона — "
                     "цена уверенной ошибки в проектировании высока",
        notes_slide_id="s20")


def build_s21(p):
    """Ландшафт 2026 — 2 колонки зрелость (cascade ≤4 названия)."""
    s = blank(p)
    slide_title(s, "Зрелое — LLM как интерфейс; генерация геометрии по тексту "
                   "— ранняя бета.", size=24, h=1.15)
    left = {
        "title": "Зрелый слой: LLM как интерфейс",
        "icon": "terminal", "icon_color": "blue",
        "lines": [
            "текстовый промпт → команды CAD над существующими "
            "детерминированными операциями",
            "под контролем инженера",
            "напр. Autodesk Assistant, Siemens Design Copilot",
        ],
    }
    right = {
        "title": "Незрелый слой: text-to-CAD / нейро-CAD",
        "icon": "flask-conical", "icon_color": "teal",
        "lines": [
            "генерация 3D-геометрии по тексту",
            "напр. Zoo.dev, Autodesk Bernini",
            "коммерческая доступность нейро-CAD заявлена «предстоящей», без "
            "даты",
        ],
    }
    two_col_compare(s, box_y=1.55, box_h=4.55, left=left, right=right)
    gold_callout(s, 0.55, 6.25, 12.25, 0.78,
                 "Сам вендор помечает Bernini «strictly experimental, not for "
                 "public use» — честный сигнал зрелости.", size=15,
                 align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """Паттерн LLM — pipeline черновик под проверку."""
    s = blank(p)
    slide_title(s, "LLM пишет черновик технического текста — инженер "
                   "верифицирует факты.", size=25)
    pipeline_row(
        s, y=1.65, h=3.30,
        stages=[
            ("Промпт с точной\nтерминологией формы",
             "требуемая структура\nи терминология документа"),
            ("LLM", "генерирует\nязыковую форму"),
            ("Черновик",
             "пояснит. записка / болванка ТУ /\nобъяснение пункта стандарта"),
            ("ОБЯЗАТЕЛЬНАЯ\nверификация инженером",
             "сверка по\nпервоисточнику"),
        ], gold_idx=3)
    gold_callout(s, 0.55, 5.30, 12.25, 1.72,
                 "LLM даёт правильную ФОРМУ и потенциально неверное "
                 "СОДЕРЖАНИЕ — полезно для черновика, опасно для финала. "
                 "Граница проходит между формой (модель) и содержанием "
                 "(инженер); верификация — всегда, без исключений.", size=16)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    """Антипаттерн — Mars Climate Orbiter dominant (cascade)."""
    s = blank(p)
    slide_title(s, "Mars Climate Orbiter: $327 млн — непроверенное "
                   "рассогласование единиц на стыке систем.", size=23, h=1.15)
    # Dominant block ~74%
    bx, by, bw, bh = 0.55, 1.45, 12.25, 4.30
    ocean_box(s, bx, by, bw, bh)
    add_image(s, icon("satellite", "blue"), bx + 0.38, by + 0.36, 1.15, 1.15)
    text_box(s, x=bx + 1.80, y=by + 0.44, w=bw - 2.1, h=0.72,
             text="Mars Climate Orbiter (1999)", size=24, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    facts = [
        ("потеря связи при выходе на орбиту Марса · стоимость миссии ",
         "≈ $327 млн"),
        ("наземный модуль выдавал фунт-силы·с, программа ожидала ньютоны·с — "
         "расхождение ×4,45", ""),
        ("конвертацию никто не выполнил: каждая сторона предположила, что её "
         "сделала другая", ""),
    ]
    fy = by + 1.55
    for pre, hi in facts:
        filled_rect(s, bx + 0.48, fy + 0.09, 0.14, 0.14, MID, radius=True,
                    radius_adj=0.5)
        runs = [{"text": pre, "size": 15.5, "color": DEEP}]
        if hi:
            runs.append({"text": hi, "size": 18, "color": GOLD, "bold": True})
        text_runs(s, bx + 0.82, fy, bw - 1.3, 0.52, runs, line_spacing=1.18)
        fy += 0.62
    gold_callout(s, bx + 0.48, by + bh - 1.10, bw - 0.96, 0.92,
                 "LLM-ассистент = такой стык систем: берёт величины из одного "
                 "контекста, передаёт в другой и УВЕРЕННО путает единицы "
                 "(psi и МПа, мм и дюймы).", size=15)
    # narrow strip — Gimli / Hyatt one line each
    sx, sy, sw, sh = 0.55, 5.92, 12.25, 1.10
    ocean_box(s, sx, sy, sw, sh)
    for j, (ic, txt) in enumerate([
        ("plane", "Gimli Glider (1983) — ручной пересчёт топлива с неверным "
         "коэффициентом"),
        ("building-2", "Hyatt Regency (1981) — изменение узла подвески «на "
         "словах» без независимого пересчёта"),
    ]):
        yy = sy + 0.16 + j * 0.42
        add_image(s, icon(ic, "blue"), sx + 0.32, yy + 0.02, 0.32, 0.32)
        text_box(s, x=sx + 0.78, y=yy, w=sw - 1.10, h=0.38, text=txt,
                 size=13, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    """Бенчмарки — range chart (без катастроф, cascade)."""
    s = blank(p)
    slide_title(s, "LLM ошибается примерно у трети–половины инженерных "
                   "расчётов.", size=25)
    subhead(s, "Воспроизводимая статистика — аргумент другого типа, чем "
               "единичные исторические катастрофы.")
    bx, by, bw, bh = 0.55, 1.80, 12.25, 2.85
    ocean_box(s, bx, by, bw, bh)
    text_box(s, x=bx + 0.35, y=by + 0.26, w=bw - 0.7, h=0.45,
             text="ORCA Benchmark — 5 SOTA-моделей на реальных расчётах",
             size=16, bold=True, color=DEEP)
    # range bar
    track_x, track_y, track_w, track_h = bx + 0.50, by + 1.15, bw - 4.2, 0.80
    filled_rect(s, track_x, track_y, track_w, track_h, SURFACE,
                stroke=LIGHT, stroke_pt=1.0, radius=True, radius_adj=0.10)
    # 45..63 of 0..100 scale
    seg_x = track_x + track_w * 0.45
    seg_w = track_w * (0.63 - 0.45)
    filled_rect(s, seg_x, track_y, seg_w, track_h, MID, radius=True,
                radius_adj=0.12)
    filled_rect(s, seg_x + seg_w - 0.07, track_y - 0.06, 0.12, track_h + 0.12,
                GOLD)
    text_box(s, x=track_x, y=track_y + track_h + 0.08, w=track_w, h=0.34,
             text="точность ~45–63%  (диапазон, не точное число)", size=13,
             italic=True, color=SLATE)
    text_box(s, x=track_x + track_w + 0.25, y=track_y - 0.05,
             w=bw - track_w - 1.1, h=0.90,
             text="лучшая ~63%", size=20, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    # two compact cards
    cw = 6.0
    gap = 0.25
    cy, ch = 4.85, 1.35
    for i, (title, body) in enumerate([
        ("EngiBench",
         "падение точности при простом перефразировании = паттерн-матчинг, "
         "не рассуждение"),
        ("Аэрокосмический кейс",
         "рекомендация по обработке нарушала 3 независимых стандарта сразу — "
         "звучала экспертно"),
    ]):
        x = 0.55 + i * (cw + gap)
        ocean_box(s, x, cy, cw, ch)
        text_box(s, x=x + 0.32, y=cy + 0.16, w=cw - 0.64, h=0.40, text=title,
                 size=16, bold=True, color=DEEP)
        text_box(s, x=x + 0.32, y=cy + 0.58, w=cw - 0.64, h=0.70, text=body,
                 size=13, color=SLATE, line_spacing=1.20)
    gold_callout(s, 0.55, 6.32, 12.25, 0.72,
                 "LLM — не источник истины по маркам/ГОСТ/допускам · "
                 "независимая верификация на каждом интерфейсе с ИИ · правка "
                 "силовой схемы = формальный change-request, не «ОК» в чате.",
                 size=14)
    speaker_notes(s, load_notes("s24"))


def build_s25(p):
    build_section_divider(
        p, here_idx=5, title="Когда сказать\nИИ «нет»",
        frame_phrase="критерии неприменимости, правильный инструмент и кто "
                     "несёт ответственность",
        notes_slide_id="s25")


def build_s26(p):
    """Критерии неприменимости — таблица 8×2 (P1: single-line ≥14pt fill≥75%
    icon-anchor)."""
    s = blank(p)
    slide_title(s, "Чем выше цена ошибки — тем меньше места вероятностному ИИ.",
                size=24, h=0.85)
    bx, by, bw, bh = 0.55, 1.28, 12.25, 5.18
    ocean_box(s, bx, by, bw, bh)
    # column headers — wider critic column, single-line cells, ≥14pt
    ic_w = 0.62
    crit_w = 3.15
    colA_w = 3.95
    colB_w = bw - 0.56 - ic_w - crit_w - colA_w
    hx0 = bx + 0.28
    hy = by + 0.16
    heads = [
        (hx0, ic_w + crit_w, "Критерий"),
        (hx0 + ic_w + crit_w, colA_w, "Почему вероятностный ИИ не годится"),
        (hx0 + ic_w + crit_w + colA_w, colB_w, "Более правильный инструмент"),
    ]
    for hx, hw, ht in heads:
        filled_rect(s, hx, hy, hw - 0.10, 0.48, MID, radius=True,
                    radius_adj=0.14)
        text_box(s, x=hx + 0.06, y=hy, w=hw - 0.22, h=0.48, text=ht,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    # Phrasings compacted to ≤4 key words per cell (P1-2 (b)) so each cell
    # is single-line legible ≥14pt @150dpi.
    rows = [
        ("file-check", "Сертификационный расчёт",
         "нужна трассируемость",
         "аттестованный решатель + верификация"),
        ("shield", "Коэффициент запаса",
         "задаётся нормой, не моделью",
         "норматив + обоснование инженера"),
        ("file-stack", "Нормоконтроль ГОСТ/ISO",
         "галлюцинирует номера и допуски",
         "нормативная база + нормоконтролёр"),
        ("scale", "Юридическая ответственность",
         "ИИ не субъект ответственности",
         "инженер с подписью"),
        ("function-square", "Есть точный решатель",
         "аппроксимация хуже точного",
         "прямой численный расчёт"),
        ("heart-pulse", "Безопасность жизни",
         "цена ошибки — жизни людей",
         "норм. расчёт + проверка + НК"),
        ("ruler", "Единицы на стыке систем",
         "уверенно путает единицы (Mars)",
         "проверка размерностей + единая СИ"),
        ("git-pull-request", "Изменение силовой схемы",
         "«по чату» без пересчёта (Hyatt)",
         "change-request + пересчёт"),
    ]
    ry0 = hy + 0.56
    rh = (bh - 0.16 - 0.56 - 0.66) / len(rows)
    for i, (ic, crit, why, tool) in enumerate(rows):
        ry = ry0 + i * rh
        band = WHITE if i % 2 == 0 else SURFACE
        filled_rect(s, bx + 0.28, ry, bw - 0.56, rh - 0.05, band,
                    stroke=COVER_OUTLINE, stroke_pt=0.75, radius=False)
        add_image(s, icon(ic, "blue"), bx + 0.38, ry + rh / 2 - 0.20,
                  0.40, 0.40)
        text_box(s, x=bx + 0.28 + ic_w, y=ry, w=crit_w - 0.05, h=rh - 0.05,
                 text=crit, size=14, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=bx + 0.28 + ic_w + crit_w, y=ry, w=colA_w - 0.14,
                 h=rh - 0.05, text=why, size=14, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, x=bx + 0.28 + ic_w + crit_w + colA_w, y=ry,
                 w=colB_w - 0.14, h=rh - 0.05, text=tool, size=14,
                 bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    gold_callout(s, bx + 0.28, by + bh - 0.66, bw - 0.56, 0.54,
                 "ИИ — инструмент расширения вариативности и черновика, не "
                 "арбитр истины и не носитель ответственности.", size=15,
                 align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """Россия нормативные границы — 4 карточки 2×2 (cascade)."""
    s = blank(p)
    slide_title(s, "В России границы ИИ-геометрии конкретны и нормативны.",
                size=25)
    cards = [
        ("file-stack", "Барьер ЕСКД", MID, [
            "выращенная органическая геометрия плохо ложится на 2D-чертёж",
            "нормоконтролируемую документацию нетривиально выпустить",
        ]),
        ("badge-check", "Аттестация и верификация", MID, [
            "ядро APM Structure3D — аттестат Ростехнадзора",
            "без аттестованной верификации в надзорное производство нельзя",
        ]),
        ("server", "187-ФЗ / 58-ФЗ", LIGHT, [
            "для значимых КИИ облачные западные CAD неприменимы",
            "блокировка AutoCAD/Fusion 2022–2024 — единая точка отказа",
        ]),
        ("copyright", "Авторство и право", GOLD, [
            "по праву РФ машина не субъект авторского права",
            "отвечает инженер или организация, не инструмент",
        ]),
    ]
    cw = 6.0
    gap = 0.25
    ch = 2.55
    gy = 0.22
    for i, (ic, title, tc, lines) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.55 + col * (cw + gap)
        y = 1.50 + row * (ch + gy)
        is_gold = (tc == GOLD)
        ocean_box(s, x, y, cw, ch,
                  fill=GOLD_TINT if is_gold else SURFACE,
                  stroke=GOLD if is_gold else LIGHT,
                  stroke_pt=2.0 if is_gold else 1.5)
        add_image(s, icon(ic, "gold" if is_gold else "blue"),
                  x + 0.34, y + 0.32, 0.70, 0.70)
        text_box(s, x=x + 1.22, y=y + 0.36, w=cw - 1.44, h=0.62, text=title,
                 size=18, bold=True, color=GOLD if is_gold else DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
        ly = y + 1.28
        for ln in lines:
            filled_rect(s, x + 0.36, ly + 0.09, 0.12, 0.12,
                        GOLD if is_gold else MID, radius=True, radius_adj=0.5)
            text_box(s, x=x + 0.60, y=ly, w=cw - 0.92, h=0.62, text=ln,
                     size=14, color=DEEP, line_spacing=1.22)
            ly += 0.66
    footer(s, "КИИ — критическая информационная инфраструктура (187-ФЗ).")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """Человек vs AI + ТРИЗ."""
    s = blank(p)
    slide_title(s, "AI генерирует форму — человек отвечает, верифицирует и "
                   "изобретает функцию.", size=24, h=1.05)
    left = {
        "title": "AI",
        "icon": "box", "icon_color": "teal",
        "lines": [
            "предлагает форму при фиксированной функции",
            "распределяет материал в заданном объёме",
        ],
    }
    right = {
        "title": "Человек",
        "icon": "pen-tool", "icon_color": "blue",
        "lines": [
            "задаёт ограничения",
            "ОТВЕЧАЕТ за результат (подпись = юр. ответственность за "
            "безопасность)",
            "ограничение деонтологическое, не техническое — не делегируется",
        ],
    }
    two_col_compare(s, box_y=1.45, box_h=3.85, left=left, right=right)
    # ТРИЗ strip
    tx, ty, tw, th = 0.55, 5.45, 12.25, 1.00
    ocean_box(s, tx, ty, tw, th)
    add_image(s, icon("lightbulb", "blue"), tx + 0.32, ty + 0.20, 0.60, 0.60)
    text_box(s, x=tx + 1.12, y=ty + 0.12, w=tw - 1.40, h=th - 0.24,
             text="ТРИЗ (Альтшуллер, с 1946) формализует ИЗОБРЕТЕНИЕ — "
                  "разрешение технического противоречия; это уровень выше, "
                  "чем оптимизация формы при фиксированной функции.",
             size=14, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.20)
    gold_callout(s, 0.55, 6.58, 12.25, 0.62,
                 "AI генерирует — человек отвечает, верифицирует и изобретает "
                 "функцию.", size=17, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    build_section_divider(
        p, here_idx=6, title="Синтез",
        frame_phrase="классы и провалы разобраны — собираем правило решения "
                     "на конкретной задаче",
        notes_slide_id="s29")


def build_s30(p):
    """Worked decision — дерево, 2 ветки × 3 вопроса."""
    s = blank(p)
    slide_title(s, "Облегчить кронштейн: один вопрос-якорь развёл "
                   "обоснованные ответы.", size=24, h=1.05)
    # task box
    tx, ty, tw, th = 3.55, 1.40, 6.25, 0.70
    filled_rect(s, tx, ty, tw, th, MID, radius=True, radius_adj=0.14)
    text_box(s, x=tx, y=ty, w=tw, h=th,
             text="Задача: снизить массу силового кронштейна", size=15,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    down_arrow(s, 6.5, ty + th + 0.02, 0.35, 0.32, fill=LIGHT)
    branches = [
        ("factory", "(а) серийная фрезеровка", [
            "Q1 нужна гарантия → детерминированный",
            "Q2 критич. ограничение — технологичность фрезы; чистая "
            "топ-оптимизация выдаст неизготовимое",
            "Q3 нужна аттестация + подпись",
        ], "Параметрика + ручная топ-оптимизация под технологичность; финал "
           "— аттестованный МКЭ. НЕ генеративный AI, НЕ суррогат для финала."),
        ("printer", "(б) аддитив", [
            "Q1 та же гарантия → детерминированный",
            "Q2 аддитив снимает ограничения формы — топ-оптимизация в полную "
            "силу",
            "Q3 финал на аттестованном решателе; суррогат — для перебора",
        ], "Детерминированная топ-оптимизация + финал на решателе; суррогат "
           "опционально для перебора; генеративный AI геометрии — нет."),
    ]
    cw = 6.0
    gap = 0.25
    cy, ch = 2.50, 3.18
    concl_h = 1.00
    for i, (ic, title, qs, concl) in enumerate(branches):
        x = 0.55 + i * (cw + gap)
        ocean_box(s, x, cy, cw, ch)
        add_image(s, icon(ic, "blue"), x + 0.28, cy + 0.20, 0.52, 0.52)
        text_box(s, x=x + 0.92, y=cy + 0.22, w=cw - 1.12, h=0.50, text=title,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        # questions area: fixed slots, sits ABOVE conclusion box
        q_area_y = cy + 0.85
        q_area_h = ch - 0.85 - concl_h - 0.20
        slot = q_area_h / len(qs)
        for j, q in enumerate(qs):
            qy = q_area_y + j * slot
            filled_rect(s, x + 0.30, qy + slot / 2 - 0.05, 0.09, 0.09, MID,
                        radius=True, radius_adj=0.5)
            text_box(s, x=x + 0.50, y=qy, w=cw - 0.78, h=slot, text=q,
                     size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.15)
        filled_rect(s, x + 0.28, cy + ch - concl_h, cw - 0.56,
                    concl_h - 0.12, GOLD_TINT, stroke=GOLD, stroke_pt=1.3,
                    radius=True, radius_adj=0.10)
        text_box(s, x=x + 0.44, y=cy + ch - concl_h, w=cw - 0.88,
                 h=concl_h - 0.12, text=concl, size=10.5, bold=True,
                 color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.16)
    gold_callout(s, 0.55, 5.85, 12.25, 0.62,
                 "Один вопрос-якорь (Q2 — что оптимизируется и какие "
                 "ограничения) развёл ответы.", size=15,
                 align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """Правило решения + полная матрица. P1: чек-лист ДОМИНИРУЕТ, матрица
    subdued."""
    s = blank(p)
    slide_title(s, "Правило решения — пять вопросов; полная матрица собрана из "
                   "всех частей.", size=23, h=1.05)
    # LEFT — checklist DOMINATES (gold box, large numbered circles, 15pt
    # bold). Slightly narrowed so the payoff matrix becomes legible, but
    # visual weight stays on the checklist (Phase-6 hierarchy preserved).
    lx, ly, lw, lh = 0.55, 1.42, 5.25, 5.03
    ocean_box(s, lx, ly, lw, lh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.2)
    add_image(s, icon("list-checks", "gold"), lx + 0.34, ly + 0.30, 0.72, 0.72)
    text_box(s, x=lx + 1.18, y=ly + 0.36, w=lw - 1.38, h=0.66,
             text="Правило решения", size=21, bold=True, color=DEEP,
             anchor=MSO_ANCHOR.MIDDLE)
    checks = [
        "Назови вид ИИ, расположи на оси",
        "Детерминированное или вероятностное?",
        "Что оптимизируется, кто задал ограничения? Полна ли постановка?",
        "Есть ли точный / нормативный инструмент?",
        "Кто отвечает?",
    ]
    chy = ly + 1.32
    for i, c in enumerate(checks):
        filled_rect(s, lx + 0.38, chy, 0.50, 0.50, MID, radius=True,
                    radius_adj=0.5)
        text_box(s, x=lx + 0.38, y=chy, w=0.50, h=0.50, text=str(i + 1),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=lx + 1.02, y=chy - 0.06, w=lw - 1.28, h=0.80, text=c,
                 size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.16)
        chy += 0.74
    # RIGHT — full 6×5 matrix (payoff). Widened + taller rows + ≥11pt
    # headers / 12pt cells so it reads from back rows (P1-5), still muted
    # (no gold box, neutral fills) so checklist keeps the focal weight.
    mx, my, mw, mh = 5.95, 1.42, 6.85, 5.03
    ocean_box(s, mx, my, mw, mh)
    add_image(s, icon("table", "blue"), mx + 0.26, my + 0.22, 0.42, 0.42)
    text_box(s, x=mx + 0.76, y=my + 0.22, w=mw - 0.95, h=0.42,
             text="Полная таксономическая матрица 6 классов", size=14,
             bold=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
    cols = ["Класс", "Природа", "Гарантирует", "Зрелость", "Кто\nотвечает"]
    cw_list = [1.78, 1.20, 1.45, 1.05, 1.07]
    tx0 = mx + 0.22
    hy = my + 0.76
    cxx = tx0
    for ci, ct in enumerate(cols):
        is_resp = (ci == 4)
        filled_rect(s, cxx, hy, cw_list[ci] - 0.06, 0.56,
                    GOLD if is_resp else LIGHT, radius=True, radius_adj=0.14)
        text_box(s, x=cxx + 0.02, y=hy, w=cw_list[ci] - 0.10, h=0.56, text=ct,
                 size=11, bold=True, color=DEEP if is_resp else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=0.95)
        cxx += cw_list[ci]
    mrows = [
        ("Топ-\nоптимизация", "детерм.", "оптимум\nфункции", "зрелая",
         "инженер"),
        ("Эволюц. / GA", "стохаст.", "оптимум\nфункции", "нишево", "инженер"),
        ("Суррогат /\nPINN", "вероятн.", "только\nв домене", "растёт",
         "инженер"),
        ("Computer\nVision", "вероятн.", "первый\nфильтр", "зрелая",
         "инженер"),
        ("LLM-\nассистент", "вероятн.", "языковую\nформу", "зрелая",
         "инженер"),
        ("Генерат.\nгеометрии", "вероятн.", "ничего\nо физике", "бета",
         "инженер"),
    ]
    ry0 = hy + 0.62
    rh = (mh - 0.76 - 0.62 - 0.50) / len(mrows)
    for ri, row in enumerate(mrows):
        ry = ry0 + ri * rh
        band = WHITE if ri % 2 == 0 else SURFACE
        filled_rect(s, tx0, ry, sum(cw_list) - 0.06, rh - 0.05, band,
                    stroke=COVER_OUTLINE, stroke_pt=0.6)
        cxx = tx0
        for ci, val in enumerate(row):
            is_resp = (ci == 4)
            text_box(s, x=cxx + 0.04, y=ry, w=cw_list[ci] - 0.10, h=rh - 0.05,
                     text=val, size=12, bold=(ci == 0 or is_resp),
                     color=GOLD if is_resp else (DEEP if ci == 0 else SLATE),
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                     line_spacing=0.98)
            cxx += cw_list[ci]
    text_box(s, x=mx + 0.22, y=my + mh - 0.48, w=mw - 0.44, h=0.40,
             text="колонка «кто отвечает» — везде инженер",
             size=11, italic=True, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    footer(s, "Профильные кафедры и суперкомпьютерные центры технических "
              "университетов работают со всем этим спектром.")
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    """Q&A — dedicated slide (Lec-N-1 pattern)."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    add_image(s, icon("message-circle-question", "gold"),
              4.55, 1.55, 1.45, 1.45)
    text_box(s, x=6.1, y=1.55, w=4.5, h=1.5, text="Q&A?",
             size=92, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)
    bx, by, bw, bh = 1.95, 3.45, 9.45, 2.55
    ocean_box(s, bx, by, bw, bh)
    text_box(s, x=bx + 0.40, y=by + 0.22, w=bw - 0.8, h=0.35,
             text="Если вопросов нет — три провокации:", size=14,
             bold=True, color=MID)
    prompts = [
        "Какой инструмент маркетинг называет «AI» в знакомом CAD — и какой "
        "это на самом деле вид ИИ?",
        "Где в проекте возникает «стык систем», на котором ИИ мог бы уверенно "
        "перепутать единицы?",
        "Назвать задачу, где ИИ обоснованно отвечают «нет», и правильный "
        "инструмент.",
    ]
    slot = (bh - 0.85) / 3
    py = by + 0.62
    for i, pr in enumerate(prompts):
        cy = py + i * slot
        filled_rect(s, bx + 0.40, cy + 0.04, 0.34, 0.34, GOLD, radius=True,
                    radius_adj=0.5)
        text_box(s, x=bx + 0.40, y=cy + 0.04, w=0.34, h=0.34,
                 text=str(i + 1), size=14, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x=bx + 0.90, y=cy, w=bw - 1.3, h=slot, text=pr,
                 size=13, italic=True, color=DEEP, anchor=MSO_ANCHOR.TOP,
                 line_spacing=1.20)
    text_box(s, x=0.55, y=6.55, w=12.3, h=0.35,
             text="Курс «AI в разных индустриях» · семинар — case study "
                  "инженерного проектирования · консультации",
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    speaker_notes(s, load_notes("s32"))


# ============================================================
# Main
# ============================================================
def main():
    p = setup_pres()
    builders = [
        build_s01, build_s02, build_s03,
        build_s04,  # divider 1
        build_s05, build_s06, build_s07, build_s08,
        build_s09,  # divider 2
        build_s10, build_s11, build_s12, build_s13, build_s14,
        build_s15,  # divider 3
        build_s16, build_s17, build_s18, build_s19,
        build_s20,  # divider 4
        build_s21, build_s22, build_s23, build_s24,
        build_s25,  # divider 5
        build_s26, build_s27, build_s28,
        build_s29,  # divider 6 (NEW — Task A)
        build_s30, build_s31, build_s32,
    ]
    for fn in builders:
        fn(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved: {OUT} ({len(builders)} slides)")


if __name__ == "__main__":
    main()
