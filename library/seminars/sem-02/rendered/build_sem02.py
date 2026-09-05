"""
Build script for Семинар 2 — «Когда ИИ, каким и на чём».

REWRITTEN for issue #182 v3 (42-slide structure, third revision, 23 owner
change requests). Replaces the stale v2 41-slide build.

=== OLD id -> NEW id mapping (v2 41-slide -> v3 42-slide) ===
This mapping exists because slide numbering SHIFTED when RAG/agent schema
slides moved and case 3 grew a beat. Do NOT assume old build_sXX function
BODY was numerically correct for the new slide of the same number -- every
slide below was rebuilt from the CURRENT slides/sNN-*.md + deck.yaml, not
from old script line numbers/comments.

  v2 id  -> v3 id  : what changed
  s01    -> s01    : same (hero cover, NASA photo kept)
  s02    -> s02    : same concept (7-case list), re-verified wording
  s03    -> s03    : divider "ИИ или обычный код?" -- NOW has bg photo + quote
  s04    -> s04    : NEW third card (целевая таблица в 1С), same 2 doc cards + photo
  s05    -> s05    : same concept (scattered formats)
  s06    -> s06    : same concept + NEW photo slot (illustration)
  s07    -> s07    : same (quickfire logs setup)
  s08    -> s08    : same (regex vs NER verdict)
  s09    -> s09    : REPLACED negative case Epic Sepsis Model -> CNET (arithmetic error)
  s10    -> s10    : divider "Встроить или своё?" -- NOW has bg photo + quote
  s11    -> s11    : same (support setup + photo)
  s12    -> s12    : same (helpdesk closed API fact)
  s13    -> s13    : same (3-option verdict)
  s14    -> s14    : same (quickfire product setup)
  s15    -> s15    : same (quickfire product verdict)
  s16    -> s16    : Kite/Copilot REBUILT as horizontal timeline (was 2-card contrast)
  (new)  -> s17    : NEW standalone slide -- Humane/Rabbit devices (was a footer note in old s16)
  s18    -> s18    : divider "Разовый вызов, RAG или агент?" -- NOW has bg photo + quote
  s17    -> s19    : RAG schema (renumbered, content ported near-verbatim)
  s18    -> s20    : Agent schema (renumbered, content ported near-verbatim)
  s20    -> s21    : case 3 setup (meeting protocol) -- ladder_row REMOVED (v3 drops all ladders)
  s22    -> s22    : case 3 verdict 1 = schema v1 (single call) -- NOW case3_schema(stage=1)
  s23    -> s23    : case 3 intro 2 (search across archive)
  s24    -> s24    : case 3 verdict 2 = schema v2 (RAG expansion) -- NOW case3_schema(stage=2)
  s25    -> s25    : case 3 intro 3 (assignments not created)
  s26    -> s26    : case 3 verdict 3 = schema v3 (agent expansion) -- NOW case3_schema(stage=3)
  s27    -> s27    : quickfire digest setup
  s28    -> s28    : quickfire digest verdict
  s29    -> s29    : REPLACED "RAG/agent in prod" content -- OLD s29 body no longer
                     applies; NEW body = Morgan Stanley (positive) + Replit (negative),
                     Octomind demoted to a footer mention only
  s28    -> s30    : divider "Внешний API или локально?" -- NOW has bg photo + quote
  s29    -> s31    : case 4 setup (sales calls) -- photo reused, renumbered
  s30    -> s32    : MERGED lawyer quote + 3-jurisdiction legal map (RU/EU/US) --
                     was a single 420-ФЗ-only card in v2
  s32    -> s33    : case 4 intro 2 (two quote cards: CTO + РОП)
  s33    -> s34    : case 4 verdict -- EXPANDED local-model detail (T-lite/T-pro,
                     4-bit quant, LoRA); calls_ladder_row REMOVED (v3 drops all ladders)
  s34    -> s35    : REPLACED JetBrains -> Apple Intelligence + Google Gemini Nano
  s35    -> s36    : open invitation ("your task") -- same concept
  s36    -> s37    : seven lessons summary -- same concept
  s37    -> s38    : homework reading -- same concept
  s38    -> s39    : reserve Klarna deep-dive -- same concept
  s39    -> s40    : reserve dedup setup -- same concept
  s40    -> s41    : reserve dedup verdict -- same concept
  s41    -> s42    : hero closing -- same concept

Ocean Gradient v3 design system. Reused helpers verbatim from the v2 script
(text_box, ocean_box, icon, chip, filled_rect, multipara_box, dashed_box,
speaker_notes, quote_block, negative_card, positive_card, footer_note,
add_image, add_image_coverfit). REMOVED: ladder_row, calls_ladder_row (v3
explicitly drops all "ход"/ladder progress strips everywhere -- the case-3
progressive schema mutation itself is the progression indicator now, see
case3_schema() below). NEW: case3_schema() shared helper for s22/s24/s26;
build_divider() extended with an optional background-photo + overlay + one-
line neutral quote param; s04 three-card layout; s16 horizontal timeline;
s32 quote+3-jurisdiction-cards layout.

Source-of-truth: deck.yaml + slides/*.md (42 slides, s01..s42; do NOT edit).
Canvas: 13.333" x 7.5" (16:9).

Gotchas (see notes/mcp-limitations.md):
  [#sem01-render-1] literal "\\n" inside a single python-pptx text run does
  not reliably line-break under LibreOffice -- always use multipara_box
  (tf.add_paragraph() per line) for multi-line text.
  [#73-render-1] add_picture(width=W, height=H) stretches non-proportionally
  -- always use add_image_coverfit (full-bleed) or add_image (one dimension
  only) for photos, never pass both width= and height= directly.
  [#sem01-render-2] python-pptx re-serializes every XML part on save() --
  raw diff is not a valid "did I only touch slide N" check.
"""
import re
import sys
from pathlib import Path

sys.path.insert(0, "/home/harness/harness-control-data/accounts/256/claude-code-klabulan-8da64c79/.local/lib/python3.12/site-packages")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED Ocean Gradient v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
NEG_TINT = RGBColor(0xFB, 0xEA, 0xEA)
NEG_LINE = RGBColor(0xB0, 0x4A, 0x4A)
POS_TINT = RGBColor(0xE9, 0xF5, 0xF2)
MUTED_FILL = RGBColor(0xEC, 0xEE, 0xF1)   # progressive-schema "already established" fill
MUTED_LINE = RGBColor(0xB9, 0xC1, 0xCB)   # progressive-schema "already established" stroke
MUTED_TEXT = RGBColor(0x8A, 0x93, 0x9E)   # progressive-schema "already established" text

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons/rendered"
SHOTS = ROOT / "assets/screenshots"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/sem-02.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Courier New"


# ============================================================
# Helpers (ported verbatim from the v2 build_sem02.py)
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")
    # python-pptx autoshapes carry a <p:style> with <a:effectRef idx="2">
    # pointing at the theme's shadow effect -- LibreOffice's PDF export applies
    # this theme effect even when spPr has an explicit empty <a:effectLst/>.
    # Removing <p:style> entirely is the reliable fix (see notes/mcp-limitations.md).
    pns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    style_el = shp._element.find(pns + "style")
    if style_el is not None:
        shp._element.remove(style_el)


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multipara_box(slide, x, y, w, h, paragraphs, *,
                   anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """Each item in `paragraphs` is a dict of text_box-style kwargs.
    Uses tf.add_paragraph() per line -- literal \\n in one run does not
    line-break reliably under LibreOffice (notes/mcp-limitations.md #sem01-render-1)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, cfg in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", align)
        p.line_spacing = cfg.get("line_spacing", 1.15)
        p.space_after = Pt(cfg.get("space_after", 0))
        r = p.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", FONT_BODY)
        r.font.size = Pt(cfg.get("size", 14))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.04, min(0.25, (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke; shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                 radius=False, radius_adj=0.16):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def dashed_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.6,
               radius_pt=12.0, dash="dash"):
    shp = ocean_box(slide, x, y, w, h, fill=fill, stroke=stroke, stroke_pt=stroke_pt,
                     radius_pt=radius_pt)
    ln = shp.line._get_or_add_ln()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    dash_el = etree.SubElement(ln, ns + "prstDash")
    dash_el.set("val", dash)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE, size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w=None, h=None):
    path = Path(path)
    if not path.exists():
        print(f"WARNING: missing image {path}")
        return None
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                        width=Inches(w), height=Inches(h))
    elif w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def icon(slide, name, color_hex, size_px, x, y, w_in):
    path = ICONS / f"{name}-{color_hex}-{size_px}.png"
    return add_image(slide, path, x, y, w=w_in, h=w_in)


def add_image_coverfit(slide, path, x, y, w, h):
    """Full-bleed / hero image helper: fills the (x,y,w,h) box exactly like
    CSS `object-fit: cover` -- no stretch distortion (see notes/mcp-limitations.md
    [#73-render-1]: passing both width= and height= to add_picture() stretches
    non-proportionally). Reads real pixel size via Pillow, sizes the picture by
    the constraining dimension, then crops the overflow off-slide via
    pic.crop_left/right/top/bottom so the box is filled edge-to-edge without
    warping the photo's aspect ratio."""
    from PIL import Image
    path = Path(path)
    if not path.exists():
        print(f"WARNING: missing image {path}")
        return None
    img_w_px, img_h_px = Image.open(path).size
    img_ratio = img_w_px / img_h_px
    box_ratio = w / h
    if img_ratio > box_ratio:
        # image wider than box -> constrain by height, crop left/right
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
        rendered_w_in = h * img_ratio
        excess_in = rendered_w_in - w
        frac = (excess_in / rendered_w_in) / 2.0
        pic.crop_left = frac
        pic.crop_right = frac
        pic.left = Inches(x)
        pic.width = Inches(w)
    else:
        # image taller than box -> constrain by width, crop top/bottom
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
        rendered_h_in = w / img_ratio
        excess_in = rendered_h_in - h
        frac = (excess_in / rendered_h_in) / 2.0
        pic.crop_top = frac
        pic.crop_bottom = frac
        pic.top = Inches(y)
        pic.height = Inches(h)
    return pic


def slide_title(slide, text, *, y=0.45, h=1.0, w=12.23, x=0.55, size=28,
                color=DEEP, bold=True, line_spacing=1.15, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=14, bold=True):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=1.2,
                radius=True, radius_adj=0.12)
    text_box(slide, x=x + 0.2, y=y + 0.06, w=w - 0.4, h=h - 0.12, text=text,
             size=size, bold=bold, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.LEFT, line_spacing=1.22)


def speaker_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def load_notes(slide_id, *, extra=None):
    """Load '## Speaker notes' verbatim from slides/{slide_id}-*.md, using the
    FIXED regex (does not truncate at '### Self-check' -- that subsection is
    INSIDE Speaker notes and must render verbatim, per
    notes/mcp-limitations.md orchestrator finding). `extra`, if given, is
    appended as one additional trailing line (used ONLY for photo attribution
    -- never edits/rewords the source text itself)."""
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r"## Speaker notes\s*\n(.*?)(?=\n## |\Z)", md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r"\n+---\s*$", "", notes)
    notes = notes.strip()
    if extra:
        notes = notes + "\n\n" + extra
    return notes


def negative_card(slide, x, y, w, h, title, body_paras, *, icon_name="x-circle"):
    filled_rect(slide, x, y, w, h, NEG_TINT, stroke=NEG_LINE, stroke_pt=1.4,
                radius=True, radius_adj=0.06)
    pad = 0.22
    icon(slide, icon_name, "21295C", 64, x + pad, y + pad, 0.4)
    text_box(slide, x + pad + 0.52, y + pad - 0.02, w - 2 * pad - 0.52, 0.45,
             text=title, size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.1)
    multipara_box(slide, x + pad, y + pad + 0.58, w - 2 * pad, h - pad - 0.6, body_paras)


def positive_card(slide, x, y, w, h, title, body_paras, *, icon_name="circle-check"):
    filled_rect(slide, x, y, w, h, POS_TINT, stroke=TEAL, stroke_pt=1.4,
                radius=True, radius_adj=0.06)
    pad = 0.22
    icon(slide, icon_name, "028090", 64, x + pad, y + pad, 0.4)
    text_box(slide, x + pad + 0.52, y + pad - 0.02, w - 2 * pad - 0.52, 0.45,
             text=title, size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.1)
    multipara_box(slide, x + pad, y + pad + 0.58, w - 2 * pad, h - pad - 0.6, body_paras)


def footer_note(slide, text, *, y=7.02):
    text_box(slide, 0.55, y, 12.23, 0.4, text=text, size=12, italic=True, color=LIGHT,
              line_spacing=1.15)


def quote_block(slide, x, y, w, h, text, *, size=15, role_icon=None):
    """Ocean-bordered quote card with a decorative quote-mark icon.
    Optional `role_icon` (Lucide icon name, e.g. "briefcase" / "scale" /
    "user-round") adds a small role-badge in the top-right corner so it's
    clear WHO is speaking."""
    ocean_box(slide, x, y, w, h, fill=SURFACE, stroke=LIGHT)
    pad = 0.24
    icon(slide, "quote", "1C7293", 64, x + pad, y + pad - 0.05, 0.32)
    text_w = w - 2 * pad - (0.5 if role_icon else 0)
    text_box(slide, x + pad, y + pad + 0.32, text_w, h - pad * 2 - 0.32,
             text=text, size=size, italic=True, color=DEEP, line_spacing=1.28,
             anchor=MSO_ANCHOR.TOP)
    if role_icon:
        badge_sz = 0.5
        bx = x + w - pad - badge_sz
        by = y + pad - 0.05
        filled_rect(slide, bx, by, badge_sz, badge_sz, SURFACE, stroke=TEAL,
                    stroke_pt=1.3, radius=True, radius_adj=0.5)
        icon(slide, role_icon, "028090", 64, bx + 0.09, by + 0.09, badge_sz - 0.18)
    return


# ============================================================
# Section-divider builder (s03 / s10 / s18 / s30) -- v3: background photo +
# overlay + neutral one-line quote. NO funnel widget, NO ladder widget
# anywhere (both removed permanently per deck.yaml revision_note_v3).
# ============================================================

def build_divider(p, sid, number, title, choice_tag, *,
                   photo_path=None, photo_credit=None, teaser=None):
    """Section divider: full-bleed background photo (muted, dark overlay for
    legibility) + big case number + title + a single choice-tag chip + an
    optional short NEUTRAL one-line comment/quote drawn from that divider's
    own speaker notes (non-spoiler -- must NOT reveal the case verdict).
    `photo_path`/`photo_credit`/`teaser` are optional so the function
    signature stays usable without a photo if one is ever unavailable."""
    s = blank(p)
    set_slide_bg(s, DEEP)
    if photo_path and Path(photo_path).exists():
        add_image_coverfit(s, photo_path, 0, 0, SLIDE_W_IN, SLIDE_H_IN)
        # dark semi-transparent overlay across the whole slide so text stays
        # readable over a photographic background (same alpha-overlay XML
        # technique as build_s01's hero photo)
        overlay = filled_rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, DEEP)
        try:
            alpha = etree.SubElement(overlay.fill.fore_color._xFill.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"),
                "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
            alpha.set("val", "78000")
        except Exception:
            pass
    # Big background number, upper-left, kept clear of the title band below it
    text_box(s, -0.1, -0.65, 6.0, 3.6, text=number, size=280, bold=True,
             color=RGBColor(0x3A, 0x44, 0x74), align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.BOTTOM, line_spacing=0.9)
    # Small gold marker dot next to the case number -- a case-index accent,
    # not decorative filler (gold-presence requirement, every slide >=1 touch)
    filled_rect(s, 2.15, 1.75, 0.22, 0.22, GOLD, radius=True, radius_adj=0.5)
    # Title, mid-left, clear of the number's descender
    text_box(s, 0.6, 3.15, 11.6, 1.5, text=title, size=40, bold=True, color=WHITE,
             line_spacing=1.1)
    # Small choice-tag chip underneath -- a label of the choice type, not a
    # spoiler of the verdict/lesson
    chip(s, 0.62, 4.75, 5.6, 0.55, choice_tag, fill=MID, size=15)
    # Optional neutral one-line comment (non-spoiler) drawn from this
    # divider's own speaker notes
    if teaser:
        text_box(s, 0.62, 5.55, 10.8, 0.7, text=teaser, size=14.5, italic=True,
                 color=RGBColor(0xC8, 0xD2, 0xDF), line_spacing=1.3)
    # photo_credit is NOT rendered on the visible slide (owner requirement:
    # zero visible photo attributions) -- it is appended to speaker notes
    # instead, via the same `extra` trailing-line mechanism used elsewhere
    # (e.g. s04/s11/s31/s41).
    extra = f"(Фото: {photo_credit})" if photo_credit else None
    speaker_notes(s, load_notes(sid, extra=extra))
    return s


# ============================================================
# Case-3 progressive schema helper (s22 / s24 / s26) -- THE single most
# important visual mechanic in the deck. Draws a FIXED set of positioned
# block-groups at the SAME x/y footprint on all 3 slides; greys out
# whatever is not yet "active" at the given stage, colors what's newly
# active. NO ladder-strip -- the schema mutation itself is the progression
# indicator (v3 explicitly removes all ladder/"ход" widgets).
# ============================================================

def _schema_block(slide, cx, cy, cw, ch, label, *, active, icon_name=None,
                   accent=MID, label_size=13):
    """One schema block. `active=False` renders muted/grey (already-
    established, de-emphasized); `active=True` renders colored (newly
    introduced at this stage)."""
    fill = SURFACE if active else MUTED_FILL
    stroke = accent if active else MUTED_LINE
    text_color = DEEP if active else MUTED_TEXT
    ocean_box(slide, cx, cy, cw, ch, fill=fill, stroke=stroke, stroke_pt=1.4 if active else 1.1)
    iy = cy + 0.14
    if icon_name:
        icon_hex = {MID: "065A82", TEAL: "028090", GOLD: "F0AB00"}.get(accent, "065A82")
        icon(slide, icon_name, icon_hex if active else "8A939E", 64, cx + (cw - 0.4) / 2, iy, 0.4)
        text_y = iy + 0.46
    else:
        text_y = cy + (ch - 0.5) / 2
    text_box(slide, cx + 0.08, text_y, cw - 0.16, ch - (text_y - cy) - 0.08,
             text=label, size=label_size, bold=active, color=text_color,
             align=PP_ALIGN.CENTER, line_spacing=1.12, anchor=MSO_ANCHOR.TOP)


def _schema_arrow(slide, x, y_mid, w, *, active):
    color = TEAL if active else MUTED_LINE
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x),
        Inches(y_mid - 0.09), Inches(w), Inches(0.18))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    disable_shadow(arrow)


def case3_schema(slide, stage, *, top=1.85):
    """Draws the case-3 pipeline at a FIXED coordinate footprint across all
    three call sites (stage 1/2/3 = s22/s24/s26). Row 1 (the v1 chain:
    письмо -> вызов модели -> протокол) is the horizontal spine, centered
    vertically in the available band. Row 0 (v2 RAG chain: хранилище ->
    индексация -> поиск -> фрагменты) sits ABOVE row 1 and feeds down into
    the "вызов модели" block. Row 2 (v3 agent chain: извлечь поручение ->
    API таск-трекера -> создать задачу -> проверить результат) sits BELOW
    row 1 and is fed FROM the "протокол" block. `top` sets the y-origin of
    row 0 so callers can reserve space above for a title/quote.

    v3 fix (iter 2): row0/row2 use SMALLER blocks + tighter gaps than v1
    (which are 4-wide vs 3-wide) so nothing overflows the 12.23in content
    width; row2's x-origin is now clamped to start at x0 (not offset right)
    so it never collides with row1 blocks; row1 is muted starting at stage
    2 (matches deck.yaml wording "приглушённое"); vertical band per row
    fixed at 1.05in with 0.55in gaps between rows -- no overlaps possible
    across any stage."""
    x0 = 0.55
    bw1, bh = 2.55, 1.0
    gap1 = 0.35

    row0_y = top
    row1_y = top + 1.55
    row2_y = top + 3.1

    # Row 1 (v1, always present from stage 1): письмо -> вызов модели -> протокол
    row1 = [("mail", "Письмо с транскриптом"), ("sparkles", "Вызов модели с шаблоном"),
            ("file-text", "Протокол")]
    row1_is_active_now = (stage == 1)
    for i, (ic, lbl) in enumerate(row1):
        cx = x0 + i * (bw1 + gap1)
        _schema_block(slide, cx, row1_y, bw1, bh, lbl,
                      active=row1_is_active_now, icon_name=ic,
                      accent=MID, label_size=12)
        if i < len(row1) - 1:
            _schema_arrow(slide, cx + bw1 + 0.04, row1_y + bh / 2, gap1 - 0.08,
                          active=row1_is_active_now)

    # Row 0 (v2 RAG chain, stage >= 2): хранилище -> индексация -> поиск -> фрагменты
    # feeds into the SAME "вызов модели" block (row1 index 1) from above.
    if stage >= 2:
        row0 = [("database", "Хранилище протоколов"), ("layers", "Индексация"),
                ("search", "Поиск по запросу"), ("file-code", "Фрагменты в модель")]
        row0_is_active_now = (stage == 2)
        bw0 = 1.72
        gap0 = 0.18
        row0_total_w = bw0 * len(row0) + gap0 * (len(row0) - 1)
        row0_x0 = x0
        for i, (ic, lbl) in enumerate(row0):
            cx = row0_x0 + i * (bw0 + gap0)
            _schema_block(slide, cx, row0_y, bw0, bh, lbl,
                          active=row0_is_active_now, icon_name=ic,
                          accent=TEAL, label_size=10.5)
            if i < len(row0) - 1:
                _schema_arrow(slide, cx + bw0 + 0.03, row0_y + bh / 2, gap0 - 0.06,
                              active=row0_is_active_now)
        # connector: last row0 block down into row1's "вызов модели" block
        connector_x = x0 + bw1 + gap1 + bw1 / 2
        conn = slide.shapes.add_connector(2, Inches(row0_x0 + row0_total_w - bw0 / 2),
                                          Inches(row0_y + bh),
                                          Inches(connector_x), Inches(row1_y))
        conn.line.color.rgb = TEAL if row0_is_active_now else MUTED_LINE
        conn.line.width = Pt(1.6)

    # Row 2 (v3 agent chain, stage >= 3): извлечь поручение -> API -> создать
    # задачу -> проверить результат -- fed FROM the row1 "протокол" block.
    # v3 fix: row2 starts at x0 (left-aligned like row0/row1), NOT offset
    # right under "протокол" -- that offset caused off-slide overflow with
    # 4 blocks. Left-aligned placement + a diagonal-reading connector from
    # "протокол" down-left into the first row2 block reads clearly enough
    # and guarantees no overflow at any stage.
    if stage >= 3:
        row2 = [("route", "Извлечь поручение"), ("cable", "Вызвать API таск-трекера"),
                ("list-checks", "Создать задачу"), ("circle-check", "Проверить результат")]
        bw2 = 1.72
        gap2 = 0.18
        row2_x0 = x0
        for i, (ic, lbl) in enumerate(row2):
            cx = row2_x0 + i * (bw2 + gap2)
            _schema_block(slide, cx, row2_y, bw2, bh, lbl,
                          active=True, icon_name=ic, accent=GOLD, label_size=10.5)
            if i < len(row2) - 1:
                _schema_arrow(slide, cx + bw2 + 0.03, row2_y + bh / 2, gap2 - 0.06,
                              active=True)
        # connector: "протокол" block (row1, index 2) down into row2 first block
        protocol_cx = x0 + (bw1 + gap1) * 2 + bw1 / 2
        conn = slide.shapes.add_connector(2, Inches(protocol_cx), Inches(row1_y + bh),
                                          Inches(row2_x0 + bw2 / 2), Inches(row2_y))
        conn.line.color.rgb = GOLD
        conn.line.width = Pt(1.6)


# ============================================================
# Slide builders -- s01..s42 (v3 42-slide case-centric structure, issue #182)
# ============================================================

def build_s01(p):
    """Hero cover -- full-bleed real photo (source: hero_cover_real_photo).
    Kept from prior sessions (same NASA/Wikimedia photo, same layout) after
    re-verification against s01-hero-cover.md wording; attempted a fresh
    6-tier search for a literal "рабочий чат/доска задач" photo this
    session (see iteration-log.md) -- no better free-licensed candidate
    found after a genuine attempt, kept the well-documented NASA photo
    (real, unstaged team-problem-solving moment) as the honest choice."""
    s = blank(p)
    set_slide_bg(s, DEEP)
    img_path = SHOTS / "s01-nasa-engineers-real.jpg"
    if img_path.exists():
        add_image_coverfit(s, img_path, 0, 0, SLIDE_W_IN, SLIDE_H_IN)
    overlay = filled_rect(s, 0, 4.6, SLIDE_W_IN, 2.9, DEEP)
    overlay.fill.fore_color.rgb = DEEP
    try:
        alpha = etree.SubElement(overlay.fill.fore_color._xFill.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"),
            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
        alpha.set("val", "80000")
    except Exception:
        pass
    text_box(s, 0.6, 4.85, 7.6, 0.45, text="СЕМИНАР 2", size=16, bold=True,
             color=GOLD, align=PP_ALIGN.LEFT)
    multipara_box(s, 0.6, 5.3, 12.1, 1.35, [
        {"text": "Когда ИИ, каким и на чём", "size": 34, "bold": True,
         "color": WHITE, "line_spacing": 1.12},
    ])
    text_box(s, 0.6, 6.35, 11.6, 0.7,
             text="Семь реальных рабочих ситуаций — и в каждой решение, которое инженер "
                  "принимает раньше, чем открывает редактор кода",
             size=15, italic=True, color=RGBColor(0xC8, 0xD2, 0xDF), line_spacing=1.28)
    speaker_notes(s, load_notes("s01",
                  extra="(Фото: NASA / Cory Huston · Wikimedia Commons · общественное достояние)"))


def build_s02(p):
    """Rebuilt per issue-182 v3 pivot review: s02-four-choices.md requires the
    MAIN block to be a plain 4-row list of the four recurring choice types
    (NO arrows, NO process schema between them), with the seven cases as a
    smaller, compact group BELOW. Previous build ignored the four-choices
    list entirely and rendered only the seven-case grid -- source mismatch,
    fixed here."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Сегодня — четыре выбора, семь кейсов", size=27, y=0.35, h=0.55)

    choices = [
        ("split", "ИИ или обычный код"),
        ("git-fork", "Встроить или делать своё"),
        ("route", "Разовый вызов, RAG или агент"),
        ("server", "Внешний API или локальный инференс"),
    ]
    choices_y = 1.05
    row_h = 0.62
    row_gap = 0.1
    for i, (ic, label) in enumerate(choices):
        cy = choices_y + i * (row_h + row_gap)
        ocean_box(s, 0.55, cy, 12.23, row_h, stroke=MID, stroke_pt=1.5)
        icon(s, ic, "065A82", 96, 0.85, cy + (row_h - 0.42) / 2, 0.42)
        text_box(s, 1.5, cy, 11.1, row_h, text=label, size=18.5,
                 bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

    cases_label_y = choices_y + 4 * (row_h + row_gap) + 0.12
    text_box(s, 0.55, cases_label_y, 12.23, 0.32,
             text="Семь кейсов, на которых сегодня разбираем эти выборы",
             size=13.5, italic=True, color=SLATE, line_spacing=1.1)

    cases = [
        ("file-text", "Документы поставщиков"),
        ("shield-alert", "Персоналка в логах"),
        ("monitor", "Помощник поддержки"),
        ("store", "Описания товаров"),
        ("list-checks", "Протоколы встреч"),
        ("mail", "Дайджест по чатам"),
        ("phone", "Звонки продаж"),
    ]
    grid_y = cases_label_y + 0.42
    row_h2 = 0.46
    gap2 = 0.1
    col_gap2 = 0.3
    col_w2 = (12.23 - col_gap2) / 2
    for i, (ic, label) in enumerate(cases):
        col = i // 4
        row = i % 4
        cx = 0.55 + col * (col_w2 + col_gap2)
        cy = grid_y + row * (row_h2 + gap2)
        ocean_box(s, cx, cy, col_w2, row_h2, stroke=LIGHT, stroke_pt=1.1)
        icon(s, ic, "065A82", 64, cx + 0.14, cy + (row_h2 - 0.3) / 2, 0.3)
        text_box(s, cx + 0.56, cy, col_w2 - 0.68, row_h2, text=label, size=12.5,
                 bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    footer_y = grid_y + 4 * (row_h2 + gap2) + 0.16
    footer_h = 7.15 - footer_y
    ocean_box(s, 0.55, footer_y, 12.23, footer_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, footer_y, 11.6, footer_h,
             text="В каждом кейсе появится деталь, которой не было в первой постановке",
             size=16.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    build_divider(p, "s03", "1", "ИИ или обычный код?", "Кейс: документы поставщиков",
                  photo_path=SHOTS / "s03-divider-code-real.jpg",
                  photo_credit="Wikimedia Commons · CC0",
                  teaser="Инженеры часто решают это на глаз, по ощущению сложности задачи")


def build_s04(p):
    """Case 1 setup -- v3: THREE cards (was 2). Two identical-format document
    cards + a NEW third card showing "целевая таблица в 1С" with the SAME
    columns, so the structural match source->target is visually obvious
    (per s04-case1-setup.md exact wording). Real photo of a document pile
    kept as a wide illustrative strip above the three cards (6-tier
    acquisition, unchanged from prior session)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Документы поставщиков", size=27, y=0.35, h=0.55)
    quote_block(s, 0.55, 0.95, 12.23, 0.95,
                "«Три человека, около 400 накладных в месяц, всё руками в 1С — можно автоматом?»",
                size=14)

    photo_y, photo_h = 2.0, 1.55
    ocean_box(s, 0.55, photo_y, 12.23, photo_h, fill=SURFACE, stroke=LIGHT)
    img_path = SHOTS / "s04-documents-pile-real.jpg"
    if img_path.exists():
        pad = 0.1
        add_image_coverfit(s, img_path, 0.55 + pad, photo_y + pad,
                            12.23 - 2 * pad, photo_h - 2 * pad)

    cards_y = photo_y + photo_h + 0.22
    ch = 2.65
    gap = 0.22
    n = 3
    cw = (12.23 - gap * (n - 1)) / n
    cols_std = ["Дата", "№ накл.", "Позиция", "Сумма"]
    cards = [
        ("A", "Документ поставщика A", cols_std, MID),
        ("Б", "Документ поставщика Б", cols_std, MID),
        ("1С", "Целевая таблица в 1С", ["Дата", "№ накладной", "Позиция", "Сумма"], GOLD),
    ]
    for i, (label, title, cols, accent) in enumerate(cards):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, cards_y, cw, ch, stroke=accent if accent == GOLD else LIGHT,
                  stroke_pt=1.8 if accent == GOLD else 1.5)
        text_box(s, cx + 0.16, cards_y + 0.12, cw - 0.32, 0.4,
                 text=title, size=12.5, bold=True, color=MID if accent != GOLD else DEEP,
                 line_spacing=1.1)
        row_h = 0.4
        table_y = cards_y + 0.58
        col_w = (cw - 0.32) / len(cols)
        for ci, colname in enumerate(cols):
            hx = cx + 0.16 + ci * col_w
            filled_rect(s, hx, table_y, col_w - 0.04, row_h,
                        accent if accent == GOLD else MID)
            text_box(s, hx, table_y, col_w - 0.04, row_h, text=colname, size=8.5,
                     bold=True, color=DEEP if accent == GOLD else WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        for r in range(2):
            ry = table_y + row_h * (r + 1) + 0.03 * (r + 1)
            for ci in range(len(cols)):
                hx = cx + 0.16 + ci * col_w
                filled_rect(s, hx, ry, col_w - 0.04, row_h, SURFACE, stroke=SOFT_GREY, stroke_pt=0.7)
                val = "24.03.2026" if (ci == 0 and r == 0) else "—"
                text_box(s, hx, ry, col_w - 0.04, row_h, text=val, size=8.5, color=SLATE,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    q_y = cards_y + ch + 0.18
    q_h = 7.5 - q_y - 0.15
    ocean_box(s, 0.55, q_y, 12.23, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, q_h, text="Что предложите?", size=20, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s04", extra="(Фото: Pizarros · Wikimedia Commons · CC BY-SA 3.0)"))


def build_s05(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Документы поставщиков", size=28)
    quote_block(s, 0.55, 1.15, 12.23, 1.0,
                "«Прислали ZIP на 4,2 ГБ — внутри и сканы актов сверки, и договоры, "
                "нужны только накладные»", size=14)

    grid_y = 2.4
    gap = 0.22
    n = 4
    cw = (12.23 - gap * (n - 1)) / n
    ch = 2.15
    variants = [
        ("A", ["Дата", "№ накл.", "Сумма"]),
        ("Б", ["Date", "Invoice#", "Total"]),
        ("В", ["Период", "Заказ", "Итого, ₽"]),
        ("Г", ["№", "от", "итого"]),
    ]
    for i, (label, cols) in enumerate(variants):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch, stroke=TEAL)
        text_box(s, cx + 0.14, grid_y + 0.12, cw - 0.28, 0.3,
                 text=f"Поставщик {label}", size=11.5, bold=True, color=MID)
        row_h = 0.36
        table_y = grid_y + 0.52
        col_w = (cw - 0.28) / len(cols)
        for ci, colname in enumerate(cols):
            hx = cx + 0.14 + ci * col_w
            filled_rect(s, hx, table_y, col_w - 0.03, row_h, TEAL)
            text_box(s, hx, table_y, col_w - 0.03, row_h, text=colname, size=8,
                     bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=0.95)
        ry = table_y + row_h + 0.03
        for ci in range(len(cols)):
            hx = cx + 0.14 + ci * col_w
            filled_rect(s, hx, ry, col_w - 0.03, row_h, SURFACE, stroke=SOFT_GREY, stroke_pt=0.7)
    text_box(s, 0.55, grid_y + ch + 0.1, 12.23, 0.35,
             text="Форматы разъехались — разные колонки, разные форматы дат, один документ на английском",
             size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

    q_y = grid_y + ch + 0.55
    ocean_box(s, 0.55, q_y, 12.23, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, 0.85, text="Меняет ли это ваш ответ?", size=22, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """v3: NEW photo slot -- a fitting illustration (someone examining
    documents / audit-style, 6-tier acquisition) alongside the verdict.
    Iter-2 fix: no free-licensed 'audit/magnifier' photo found after a
    genuine search this session (see iteration-log.md) -- honest icon-based
    fallback, enlarged to fill the vertical band properly (was a small
    centered icon floating in a mostly-empty box).
    Iter-3 fix (issue-182 pivot review round 2): removed the footer_note
    entirely -- "Держите в голове риск из Лекции 1: модель тихо деградирует
    на форматах..." is not present in s06-case1-verdict.md and owner
    explicitly asked to drop it (No Extra Content Rule)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Документы поставщиков — разбор", size=25, y=0.35, h=0.55)

    verdict_y = 1.05
    verdict_h = 2.9
    photo_w = 3.3
    ocean_box(s, 0.55, verdict_y, 12.23 - photo_w - 0.25, verdict_h, fill=SURFACE,
              stroke=MID, stroke_pt=1.6)
    icon(s, "workflow", "065A82", 96, 0.85, verdict_y + 0.25, 0.55)
    multipara_box(s, 1.65, verdict_y + 0.22, 12.23 - photo_w - 0.25 - 1.2, verdict_h - 0.44, [
        {"text": "Гибрид: LLM-извлечение по вариативному входу + жёсткая валидация "
                 "кодом на выходе", "size": 16.5, "bold": True, "color": DEEP,
         "line_spacing": 1.25, "space_after": 10},
        {"text": "Документ (любой формат) → LLM извлекает поля → код проверяет "
                 "типы/обязательные значения/диапазоны → данные в 1С", "size": 13,
         "color": SLATE, "line_spacing": 1.32, "space_after": 10},
        {"text": "Разумный ответ здесь — гибрид, а не выбор одного из двух подходов",
         "size": 12, "italic": True, "color": SLATE, "line_spacing": 1.28},
    ])
    photo_x = 12.23 + 0.55 - photo_w
    ocean_box(s, photo_x, verdict_y, photo_w, verdict_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.5)
    img_path = SHOTS / "s06-audit-magnifier-real.jpg"
    if img_path.exists():
        pad = 0.1
        add_image_coverfit(s, img_path, photo_x + pad, verdict_y + pad,
                            photo_w - 2 * pad, verdict_h - 2 * pad)
    else:
        icon(s, "search", "065A82", 96, photo_x + (photo_w - 1.4) / 2, verdict_y + 0.5, 1.4)
        text_box(s, photo_x + 0.2, verdict_y + 2.15, photo_w - 0.4, 0.6,
                 text="изучайте данные, не верьте на слово", size=11.5, italic=True,
                 color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.25)

    lesson_y = verdict_y + verdict_h + 0.3
    ocean_box(s, 0.55, lesson_y, 12.23, 1.3, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, lesson_y, 11.6, 1.3,
             text="Два примера — не выборка. Просите представление полного объёма данных",
             size=19, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    speaker_notes(s, load_notes("s06",
                  extra="(Фото: Wikimedia Commons · CC BY-SA)" if img_path.exists() else None))


def build_s07(p):
    """Iter-2 fix: scene + question boxes enlarged to close ~1.6in of dead
    space at the bottom of the slide (visual mass balance)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Персоналка в логах", size=28)
    quote_block(s, 0.55, 1.15, 12.23, 1.05,
                "«40 ГБ логов за квартал — подрядчик ждёт архив в пятницу»",
                size=16)

    scene_y = 2.5
    scene_h = 2.9
    ocean_box(s, 0.55, scene_y, 12.23, scene_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "file-text", "065A82", 96, 0.95, scene_y + 0.55, 0.95)
    icon(s, "phone", "028090", 64, 2.25, scene_y + 0.8, 0.55)
    icon(s, "mail", "065A82", 96, 3.1, scene_y + 0.75, 0.6)
    icon(s, "user-round", "F0AB00", 64, 3.95, scene_y + 0.8, 0.55)
    multipara_box(s, 4.9, scene_y + 0.6, 7.0, 1.9, [
        {"text": "Одна строка лога — телефон, email и фраза «передайте это Ивану "
                 "из бухгалтерии»", "size": 17, "bold": True, "color": DEEP,
         "line_spacing": 1.3, "space_after": 10},
        {"text": "Подрядчику нужны структура запросов, тайминги, ошибки — не данные людей",
         "size": 14, "color": SLATE, "line_spacing": 1.3},
    ])

    q_y = scene_y + scene_h + 0.3
    q_h = 7.15 - q_y
    ocean_box(s, 0.55, q_y, 12.23, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, q_h, text="Как максимально быстро и дёшево очистить?",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """Iter-2 fix: typo 'detекторская' (mixed-script garbage) replaced with
    correct Russian wording; card bottom padding filled (was ~30% empty
    space) with an explanatory sentence instead of a short chip."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Персоналка в логах — разбор", size=26)

    grid_y = 1.5
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 3.55
    ocean_box(s, 0.55, grid_y, cw, ch, stroke=TEAL)
    icon(s, "hash", "028090", 96, 0.55 + 0.25, grid_y + 0.22, 0.55)
    text_box(s, 0.55 + 0.95, grid_y + 0.22, cw - 1.2, 0.55, text="Regex", size=17,
             bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 0.55 + 0.25, grid_y + 1.0, cw - 0.5, 1.15,
             text="Телефоны, email — чёткий формат, дёшево и надёжно вычищает обычный regex",
             size=14.5, color=DEEP, line_spacing=1.3)
    chip(s, 0.55 + 0.25, grid_y + 2.15, cw - 0.5, 0.42,
         "Детерминированная логика — модель не нужна", fill=SURFACE, stroke=TEAL,
         color=MID, size=11)
    text_box(s, 0.55 + 0.25, grid_y + 2.75, cw - 0.5, 0.7,
             text="Никакая модель здесь не нужна и только замедлила бы обработку 40 ГБ логов",
             size=11.5, italic=True, color=SLATE, line_spacing=1.25)

    rx = 0.55 + cw + gap
    ocean_box(s, rx, grid_y, cw, ch, stroke=MID)
    icon(s, "search", "065A82", 96, rx + 0.25, grid_y + 0.22, 0.55)
    text_box(s, rx + 0.95, grid_y + 0.22, cw - 1.2, 0.55, text="NER / малая модель", size=16,
             bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.25, grid_y + 1.0, cw - 0.5, 1.15,
             text="Имена людей, косвенные упоминания в свободном тексте — нужен контекст",
             size=14, color=DEEP, line_spacing=1.28)
    text_box(s, rx + 0.25, grid_y + 2.1, cw - 0.5, 0.42,
             text="slovnet (Natasha) · ~30 МБ · CPU · ~25 статей/сек",
             size=11.5, italic=True, color=SLATE, line_spacing=1.2)
    text_box(s, rx + 0.25, grid_y + 2.7, cw - 0.5, 0.75,
             text="~60× меньше BERT сравнимого класса, потеря качества 1-2 п.п.",
             size=11.5, italic=True, color=SLATE, line_spacing=1.25)

    footer_y = grid_y + ch + 0.25
    ocean_box(s, 0.55, footer_y, 12.23, 1.0, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, footer_y, 11.6, 1.0,
             text="Граница между regex и моделью проходит внутри одной задачи, не между задачами",
             size=16, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """v3: REPLACED negative case Epic Sepsis Model -> CNET (arithmetic
    error). Positive case Ramp unchanged."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Когда нужна была формула — и когда ИИ закрыл нерешаемое", size=22)

    grid_y = 1.5
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 2.95
    negative_card(s, 0.55, grid_y, cw, ch, "CNET, январь 2023", [
        {"text": "Из 77 ИИ-сгенерированных статей о финансах исправления потребовались 41",
         "size": 13, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "«10 000 $ под 3% годовых → 10 300 $ за год» вместо корректных 300 $ дохода",
         "size": 12.5, "bold": True, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Расчёт должен делать детерминированная формула, а не LLM «на глаз»",
         "size": 12.5, "italic": True, "color": DEEP, "line_spacing": 1.25},
    ], icon_name="x-circle")

    rx = 0.55 + cw + gap
    positive_card(s, rx, grid_y, cw, ch, "Ramp, merchant classification", [
        {"text": "Вручную: покрытие заявок 1,5-3% → после LLM-агента почти полное",
         "size": 13, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Обработка < 10 сек, ~99% решений корректны", "size": 12.5, "bold": True,
         "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "~25% заявок агент обоснованно отклоняет", "size": 12.5,
         "color": DEEP, "line_spacing": 1.25},
    ], icon_name="circle-check")

    takeaway_y = grid_y + ch + 0.28
    ocean_box(s, 0.55, takeaway_y, 12.23, 1.15, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, takeaway_y, 11.6, 1.15,
             text="Ramp — замена ручной обработки с мизерным покрытием, не замена работавшего "
                  "regex-решения",
             size=17, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    build_divider(p, "s10", "2", "Встроить или делать своё?", "Кейс: помощник поддержки",
                  photo_path=SHOTS / "s10-divider-workspace-real.jpg",
                  photo_credit="Oliver Propst · Wikimedia Commons · CC BY-SA 3.0",
                  teaser="Скрытая цена почти никогда не видна в демо на презентации идеи")


def build_s11(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Помощник поддержки", size=28)
    quote_block(s, 0.55, 1.1, 12.23, 1.35,
                "«Сделайте нам ИИ-помощника для операторов — подсказки по тикетам, "
                "поиск похожих обращений»", size=16)

    photo_x, photo_y, photo_w, photo_h = 0.55, 2.65, 4.3, 3.95
    ocean_box(s, photo_x, photo_y, photo_w, photo_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    img_path = SHOTS / "s11-support-headset-real.jpg"
    if img_path.exists():
        pad = 0.12
        add_image_coverfit(s, img_path, photo_x + pad, photo_y + pad,
                            photo_w - 2 * pad, photo_h - 2 * pad)

    rx = photo_x + photo_w + 0.3
    rw = 12.23 - photo_w - 0.3
    ocean_box(s, rx, photo_y, rw, 1.65, fill=SURFACE, stroke=LIGHT)
    text_box(s, rx + 0.2, photo_y + 0.16, rw - 0.4, 1.35,
             text="12 операторов на линии, около 300 тикетов в день, пик 10:00-13:00",
             size=15.5, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    q_y = photo_y + 1.95
    q_h = photo_h - 1.95
    ocean_box(s, rx, q_y, rw, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, rx + 0.2, q_y, rw - 0.4, q_h, text="Где должен жить этот помощник?",
             size=24, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    speaker_notes(s, load_notes("s11",
                  extra="(Фото: FiveOne51 · Wikimedia Commons · CC BY-SA 3.0)"))


def build_s12(p):
    """Iter-2 fix: fact box was 3.4in tall with ~1in of actual content
    vertically centered (large dead margin top/bottom) -- shrunk box,
    enlarged icon + text, and added a supporting detail line to use the
    freed space productively instead of leaving it centered in a void."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Помощник поддержки", size=28)
    quote_block(s, 0.55, 1.05, 12.23, 1.15,
                "«Открытого API нет, кастомизаций нет — вендор и на письма не отвечает»",
                size=15, role_icon="briefcase")

    q_y = 2.45
    ocean_box(s, 0.55, q_y, 12.23, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, 0.85, text="Что делаем?", size=24, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)

    fact_y = q_y + 1.1
    fact_h = 3.4
    ocean_box(s, 0.55, fact_y, 12.23, fact_h, fill=SURFACE, stroke=LIGHT)
    icon(s, "lock", "21295C", 96, 0.85, fact_y + 0.35, 0.85)
    multipara_box(s, 1.95, fact_y + 0.35, 10.4, fact_h - 0.6, [
        {"text": "Коробочный helpdesk — систему купили 4 года назад", "size": 20, "bold": True,
         "color": DEEP, "line_spacing": 1.28, "space_after": 12},
        {"text": "На ней же тикеты HR и фасилити, контракт с вендором ещё на 2 года — "
                 "нет открытого API, нет кастомизаций, вендор не отвечает даже на письма",
         "size": 15, "color": SLATE, "line_spacing": 1.35, "space_after": 12},
        {"text": "Компания написала вендору с вопросом об интеграции — ответа так и не получила",
         "size": 13, "italic": True, "color": SLATE, "line_spacing": 1.3},
    ])
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Помощник поддержки — разбор", size=26)

    grid_y = 1.4
    gap = 0.28
    cw = (12.23 - gap * 2) / 3
    ch = 3.9
    options = [
        ("split", "Решение сбоку", "Расширение браузера или вторая панель",
         "Два места одновременно, риск: оператор не будет пользоваться"),
        ("repeat", "Смена вендора", "Helpdesk-система с открытым API",
         "Миграция данных, переобучение команды, риск простоя"),
        ("route", "Пересмотр задачи", "Асинхронный помощник, готовит черновики заранее",
         "Оператор открывает готовое там, где удобно, но нужен отдельный процесс"),
    ]
    for i, (ic, title, desc, cost) in enumerate(options):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, grid_y, cw, ch)
        icon(s, ic, "065A82", 96, cx + (cw - 0.55) / 2, grid_y + 0.24, 0.55)
        text_box(s, cx + 0.16, grid_y + 0.95, cw - 0.32, 0.5, text=title, size=15,
                 bold=True, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.1)
        text_box(s, cx + 0.16, grid_y + 1.5, cw - 0.32, 1.1, text=desc, size=12,
                 color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.25)
        filled_rect(s, cx + 0.14, grid_y + ch - 1.35, cw - 0.28, 1.15, GOLD_TINT,
                    stroke=GOLD, stroke_pt=1.1, radius=True, radius_adj=0.1)
        text_box(s, cx + 0.24, grid_y + ch - 1.28, cw - 0.48, 1.0, text=f"Цена: {cost}",
                 size=10.5, italic=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.2,
                 anchor=MSO_ANCHOR.MIDDLE)

    footer_y = grid_y + ch + 0.2
    ocean_box(s, 0.55, footer_y, 12.23, 0.6, fill=SURFACE, stroke=LIGHT)
    text_box(s, 0.85, footer_y, 11.6, 0.6, text="Единственно верного ответа нет",
             size=15, bold=True, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """Iter-2 fix: scene + question boxes enlarged (was ~1.8in dead space
    at bottom)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Описания товаров", size=28)
    quote_block(s, 0.55, 1.15, 12.23, 1.1,
                "«Маркетинг просит отдельный ИИ-сервис для генерации описаний товаров»",
                size=16.5)

    scene_y = 2.55
    scene_h = 2.85
    ocean_box(s, 0.55, scene_y, 12.23, scene_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "store", "065A82", 96, 0.95, scene_y + (scene_h - 1.1) / 2, 1.1)
    icon(s, "arrow-right", "F0AB00", 64, 2.35, scene_y + (scene_h - 0.6) / 2, 0.6)
    icon(s, "sparkles", "065A82", 96, 3.25, scene_y + (scene_h - 1.1) / 2, 1.1)
    multipara_box(s, 4.7, scene_y + (scene_h - 1.6) / 2, 7.1, 1.6, [
        {"text": "12 000 товарных позиций, сейчас описания пишет один копирайтер на аутсорсе",
         "size": 16.5, "bold": True, "color": DEEP, "line_spacing": 1.3, "space_after": 10},
        {"text": "Новый отдельный инструмент: логин, вкладка, форма ввода характеристик товара",
         "size": 14, "color": SLATE, "line_spacing": 1.3},
    ])

    q_y = scene_y + scene_h + 0.3
    q_h = 7.15 - q_y
    ocean_box(s, 0.55, q_y, 12.23, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, q_h, text="Отдельный сервис — хорошая идея?",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Описания товаров — разбор", size=27)

    row_y = 1.55
    row_h = 2.0
    ocean_box(s, 0.55, row_y, 12.23, row_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "store", "065A82", 96, 0.9, row_y + (row_h - 0.85) / 2, 0.85)
    icon(s, "arrow-right", "F0AB00", 64, 2.15, row_y + (row_h - 0.5) / 2, 0.5)
    icon(s, "zap", "F0AB00", 96, 3.05, row_y + (row_h - 0.85) / 2, 0.85)
    multipara_box(s, 4.2, row_y, 7.6, row_h, [
        {"text": "Кнопка «сгенерировать описание»", "size": 20, "bold": True, "color": DEEP,
         "line_spacing": 1.2, "space_after": 6},
        {"text": "в существующей админке каталога, которой маркетинг и так пользуется каждый день",
         "size": 15, "color": DEEP, "line_spacing": 1.3},
    ], anchor=MSO_ANCHOR.MIDDLE)

    fn_y = row_y + row_h + 0.35
    ocean_box(s, 0.55, fn_y, 12.23, 1.2, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, fn_y, 11.6, 1.2,
             text="Отдельный продукт — второй логин и вкладка, цена, о которой не подумали",
             size=17, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    speaker_notes(s, load_notes("s15"))


def build_s16(p):
    """v3: REBUILT as a horizontal timeline with 5 dated points (was a
    2-card contrast in v2). Per s16-kite-vs-copilot.md: 2014 Kite founded ->
    2019 local inference -> 29.06.2021 Copilot technical preview ->
    21.06.2022 Copilot GA $10/mo -> 21.11.2022 Kite shut down. Highlighted
    "~17 месяцев пересечения" span between technical-preview and shutdown.
    User-count callouts below."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Kite и Copilot", size=27, y=0.35, h=0.55)
    text_box(s, 0.55, 0.95, 12.23, 0.4,
             text="Одна технология, разная встроенность — 17 месяцев прямой конкуренции",
             size=14, italic=True, color=MID)

    track_y = 3.05
    track_x0, track_x1 = 0.75, 12.6
    track = filled_rect(s, track_x0, track_y - 0.02, track_x1 - track_x0, 0.05, SOFT_GREY)

    events = [
        (2014.0, "2014", "Kite основан", False),
        (2019.0, "2019", "Локальное исполнение модели", False),
        (2021.5, "29.06.2021", "Copilot technical preview", False),
        (2022.47, "21.06.2022", "Copilot GA, $10/мес", False),
        (2022.9, "21.11.2022", "Kite закрыт", True),
    ]
    dmin, dmax = 2013.6, 2023.3

    def to_x(d):
        return track_x0 + (d - dmin) / (dmax - dmin) * (track_x1 - track_x0)

    # highlighted overlap span between technical-preview (idx 2) and Kite
    # shutdown (idx 4) -- iter-2 fix: moved WELL above the track (was
    # colliding with the "above" date labels at track_y-1.15) and widened
    # so the 1-line caption never wraps to 2-3 lines
    ox0, ox1 = to_x(events[2][0]), to_x(events[4][0])
    filled_rect(s, ox0, track_y - 0.02, ox1 - ox0, 0.05, GOLD)
    span_label_w = max(ox1 - ox0 + 1.4, 3.6)
    text_box(s, (ox0 + ox1) / 2 - span_label_w / 2, track_y - 1.85, span_label_w, 0.3,
             text="~17 месяцев пересечения на рынке",
             size=12, bold=True, italic=True, color=DEEP, align=PP_ALIGN.CENTER)

    for i, (d, date_lbl, ev_lbl, is_pivot) in enumerate(events):
        x = to_x(d)
        r = 0.09 if not is_pivot else 0.13
        dotcolor = GOLD if is_pivot else MID
        filled_rect(s, x - r, track_y - r + 0.015, r * 2, r * 2, dotcolor, radius=True, radius_adj=0.5)
        above = (i % 2 == 0)
        label_y = track_y - 1.4 if above else track_y + 0.25
        text_box(s, x - 1.05, label_y, 2.1, 0.3, text=date_lbl, size=12,
                 bold=True, color=DEEP if not is_pivot else GOLD, align=PP_ALIGN.CENTER)
        text_box(s, x - 1.15, label_y + (0.3 if above else 0.0) + (0.0 if above else 0.32),
                 2.3, 0.55, text=ev_lbl, size=10.5, color=SLATE, align=PP_ALIGN.CENTER,
                 line_spacing=1.15)

    users_y = 4.85
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 1.15
    ocean_box(s, 0.55, users_y, cw, ch, fill=NEG_TINT, stroke=NEG_LINE, stroke_pt=1.3)
    text_box(s, 0.55 + 0.2, users_y + 0.1, cw - 0.4, ch - 0.2,
             text="Kite — ~500 000 пользователей на пике", size=13.5, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    rx = 0.55 + cw + gap
    ocean_box(s, rx, users_y, cw, ch, fill=POS_TINT, stroke=TEAL, stroke_pt=1.3)
    text_box(s, rx + 0.2, users_y + 0.1, cw - 0.4, ch - 0.2,
             text="Copilot — ~20 млн пользователей · 4,7 млн платных · ~90% Fortune 100",
             size=12.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    takeaway_y = users_y + ch + 0.22
    ocean_box(s, 0.55, takeaway_y, 12.23, 0.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, takeaway_y, 11.6, 0.95,
             text="Встроенность в уже существующий рабочий поток победила отдельный продукт "
                  "при технологии сравнимого класса",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """NEW standalone slide (v3): standalone devices, Humane AI Pin + Rabbit
    R1. Humane: NO free-licensed photo found after a genuine Commons search
    this session (opensearch + category browse + keyword search all
    returned zero relevant hits, confirming the brief's own 404 finding for
    File:Humane-press-aipin-family.png) -- built as an icon-based device
    card, per brief instructions, honestly logged in iteration-log.md.
    Rabbit R1: real CC BY 3.0 photo found (Booredatwork.com via Commons)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Отдельное устройство против фичи в телефоне", size=23, y=0.35, h=0.55)

    grid_y = 1.15
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 4.15
    # Humane AI Pin -- icon-based card (no free photo found, honest fallback)
    ocean_box(s, 0.55, grid_y, cw, ch, fill=NEG_TINT, stroke=NEG_LINE, stroke_pt=1.4)
    icon(s, "smartphone", "21295C", 64, 0.55 + 0.24, grid_y + 0.2, 0.5)
    text_box(s, 0.55 + 0.9, grid_y + 0.2, cw - 1.1, 0.5, text="Humane AI Pin", size=16,
             bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    multipara_box(s, 0.55 + 0.24, grid_y + 0.95, cw - 0.48, ch - 1.1, [
        {"text": "Носимый значок без экрана, лазерный проектор на ладонь, голосовое "
                 "управление · $699 + подписка $24/мес",
         "size": 12.5, "color": DEEP, "line_spacing": 1.28, "space_after": 8},
        {"text": "Продажи с апреля 2024", "size": 12, "color": SLATE, "line_spacing": 1.2,
         "space_after": 6},
        {"text": "Привлечено ~$241 млн → активы проданы HP за $116 млн (февраль 2025) → "
                 "серверы отключены 28.02.2025",
         "size": 12.5, "bold": True, "color": DEEP, "line_spacing": 1.28},
    ])

    rx = 0.55 + cw + gap
    ocean_box(s, rx, grid_y, cw, ch, fill=SURFACE, stroke=TEAL, stroke_pt=1.5)
    img_path = SHOTS / "s17-rabbit-r1-real.jpg"
    photo_h = 2.15
    if img_path.exists():
        pad = 0.14
        add_image_coverfit(s, img_path, rx + pad, grid_y + pad, cw - 2 * pad, photo_h - pad)
    text_box(s, rx + 0.24, grid_y + photo_h + 0.14, cw - 0.48, 0.4, text="Rabbit R1", size=16,
             bold=True, color=DEEP)
    text_box(s, rx + 0.24, grid_y + photo_h + 0.56, cw - 0.48, ch - photo_h - 0.7,
             text="Продано ~100 000 устройств, по оценкам активно пользуются около 5% купивших",
             size=12.5, color=DEEP, line_spacing=1.28)

    takeaway_y = grid_y + ch + 0.22
    ocean_box(s, 0.55, takeaway_y, 12.23, 1.1, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, takeaway_y, 11.6, 1.1,
             text="Отдельное устройство конкурирует не только с технологией, но и с "
                  "привычкой не переключаться, когда функция уже есть в телефоне",
             size=14.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.22)
    speaker_notes(s, load_notes("s17",
                  extra="(Фото Rabbit R1: Booredatwork.com · Wikimedia Commons · CC BY 3.0. "
                        "Фото Humane AI Pin не найдено — свободно лицензированного изображения "
                        "на Wikimedia Commons не существует, проверено через API-поиск и обзор "
                        "категорий; см. iteration-log.md.)"))


def build_s18(p):
    build_divider(p, "s18", "3", "Разовый вызов, RAG или агент?", "Кейс: протоколы встреч",
                  photo_path=SHOTS / "s18-divider-serverroom-real.jpg",
                  photo_credit="SimonWaldherr · Wikimedia Commons · CC BY-SA 4.0",
                  teaser="Чем внушительнее задача, тем сложнее архитектура — интуицию стоит проверять")


def build_s19(p):
    """RAG schema -- standalone slide, positioned BEFORE case 3, ported
    near-verbatim from the prior session's build."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Два слова про RAG", size=28)
    text_box(s, 0.55, 1.05, 9.5, 0.5,
             text="Поиск по базе знаний подкладывает нужные фрагменты в контекст модели",
             size=15, italic=True, color=MID, line_spacing=1.2)
    chip(s, 11.15, 1.0, 1.65, 0.42, "Лекция 3 →", fill=TEAL, size=11.5)

    steps = [
        ("search", "Вопрос пользователя"),
        ("database", "Поиск фрагментов\nв базе знаний"),
        ("layers", "Фрагменты →\nв контекст модели"),
        ("sparkles", "Модель генерирует\nответ"),
    ]
    grid_y = 2.35
    grid_h = 2.9
    n = 4
    gap_arrow = 0.55
    cw = (12.23 - gap_arrow * (n - 1)) / n
    for i, (ic, lbl) in enumerate(steps):
        cx = 0.55 + i * (cw + gap_arrow)
        ocean_box(s, cx, grid_y, cw, grid_h)
        icon(s, ic, "065A82", 96, cx + (cw - 0.6) / 2, grid_y + 0.35, 0.6)
        lines = lbl.split("\n")
        paras = [{"text": ln, "size": 13.5, "bold": True, "color": DEEP,
                  "align": PP_ALIGN.CENTER, "line_spacing": 1.15} for ln in lines]
        multipara_box(s, cx + 0.12, grid_y + 1.15, cw - 0.24, 1.6, paras, align=PP_ALIGN.CENTER)
        chip(s, cx + 0.14, grid_y + grid_h - 0.5, 0.34, 0.34, str(i + 1), fill=GOLD,
             color=DEEP, size=13)
        if i < n - 1:
            ax = cx + cw + 0.06
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax),
                Inches(grid_y + grid_h / 2 - 0.16), Inches(gap_arrow - 0.12), Inches(0.32))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()
            disable_shadow(arrow)

    ocean_box(s, 0.55, 5.55, 12.23, 1.05, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, 5.55, 11.6, 1.05,
             text="Когда нужен: база большая, меняется, весь объём не помещается в один запрос",
             size=16.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    """Agent schema -- standalone slide, positioned BEFORE case 3, ported
    near-verbatim from the prior session's build."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Два слова про агента", size=28)
    text_box(s, 0.55, 1.05, 9.5, 0.5,
             text="Модель работает циклом с внешними инструментами, а не одним вызовом",
             size=15, italic=True, color=MID, line_spacing=1.2)
    chip(s, 11.15, 1.0, 1.65, 0.42, "Лекция 3 →", fill=TEAL, size=11.5)

    cx0, r_w = 6.67, 4.7
    positions = [
        (cx0 - r_w / 2 - 1.55, 2.0, "git-fork", "Строит план"),
        (cx0 + r_w / 2 - 1.55, 2.0, "wrench", "Вызывает внешний\nинструмент (API,\nпоиск, код)"),
        (cx0 + r_w / 2 - 1.55, 3.9, "circle-check", "Проверяет\nрезультат вызова"),
        (cx0 - r_w / 2 - 1.55, 3.9, "route", "Решает, что\nделать дальше"),
    ]
    card_w, card_h = 3.1, 1.6
    for i, (cx, cy, ic, lbl) in enumerate(positions):
        ocean_box(s, cx, cy, card_w, card_h)
        icon(s, ic, "065A82", 96, cx + 0.16, cy + 0.16, 0.44)
        chip(s, cx + card_w - 0.48, cy + 0.14, 0.32, 0.32, str(i + 1), fill=GOLD,
             color=DEEP, size=11.5)
        lines = lbl.split("\n")
        paras = [{"text": ln, "size": 12, "bold": True, "color": DEEP,
                  "line_spacing": 1.1} for ln in lines]
        multipara_box(s, cx + 0.16, cy + 0.68, card_w - 0.32, card_h - 0.78, paras)

    loop_sz = 1.0
    lcx = cx0 - loop_sz / 2
    lcy = 2.0 + card_h + 0.05
    filled_rect(s, lcx, lcy, loop_sz, loop_sz, TEAL, radius=True, radius_adj=0.5)
    icon(s, "repeat", "FFFFFF", 48, lcx + 0.2, lcy + 0.2, 0.6)
    text_box(s, cx0 - 1.1, lcy + loop_sz + 0.04, 2.2, 0.3, text="цикл повторяется",
             size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

    ocean_box(s, 0.55, 5.95, 12.23, 1.05, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, 5.95, 11.6, 1.05,
             text="Когда нужен: несколько шагов и действия во внешних системах",
             size=16.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s20"))


def build_s21(p):
    """Case 3 setup -- v3: NO ladder_row (v3 drops all ladder/ход strips
    everywhere; the case3_schema mutation on s22/s24/s26 IS the progression
    indicator now). Iter-2 fix: photo + right column enlarged to close
    ~1.75in of dead space below the photo row."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Протоколы встреч", size=28)
    quote_block(s, 0.55, 1.1, 12.23, 1.2,
                "«После каждого созвона транскрипт приходит письмом, нужен протокол: "
                "решения, поручения, сроки»", size=15)

    photo_x, photo_y, photo_w, photo_h = 0.55, 2.55, 6.9, 4.6
    ocean_box(s, photo_x, photo_y, photo_w, photo_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    img_path = SHOTS / "s20-meeting-room-real.jpg"
    if img_path.exists():
        pad = 0.12
        add_image_coverfit(s, img_path, photo_x + pad, photo_y + pad,
                            photo_w - 2 * pad, photo_h - 2 * pad)

    rx = photo_x + photo_w + 0.3
    rw = 12.23 - photo_w - 0.3
    ocean_box(s, rx, photo_y, rw, 2.55, fill=SURFACE, stroke=LIGHT)
    text_box(s, rx + 0.2, photo_y + 0.18, rw - 0.4, 2.19,
             text="Около 15 встреч в неделю, протокол сейчас пишет проджект вручную",
             size=15.5, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    q_y = photo_y + 2.85
    q_h = photo_h - 2.85
    ocean_box(s, rx, q_y, rw, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, rx + 0.2, q_y, rw - 0.4, q_h, text="Какая архитектура нужна?",
             size=20, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    speaker_notes(s, load_notes("s21",
                  extra="(Фото: Amtec Photos · Wikimedia Commons · CC BY 2.0)"))


def build_s22(p):
    """Case 3 verdict 1 = schema v1 (single call). v3: uses case3_schema
    (stage=1) instead of a plain 3-block row -- NO ladder_row.

    IMPORTANT iter-2 fix: case3_schema keeps row1 (the v1 chain) at a FIXED
    absolute y-offset from `top` (top+1.55) regardless of stage, BY DESIGN
    -- this is what makes "the same schema growing" read visually identical
    across s22/s24/s26 (the whole point of the mechanic). The iter-1 bug was
    in THIS call site: it tried to compress the schema into a short frame
    box sized for only 1 visible row, which pushed row1 off the bottom of
    its own frame. Fixed by giving s22 the SAME frame height/top as s24/s26
    use for their full 3-row footprint -- row1 sits in the correct fixed
    slot with generous matching whitespace above/below it at stage 1 (that
    whitespace becomes row0/row2 at stages 2/3, so this is consistent
    across the trilogy, not a defect)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Протоколы встреч — разбор", size=22, y=0.2, h=0.4)
    quote_block(s, 0.55, 0.62, 12.23, 0.95,
                "«Один вызов, ни хранилища, ни цикла»", size=14)

    schema_y, schema_h = 1.68, 4.3
    ocean_box(s, 0.55, schema_y, 12.23, schema_h, fill=SURFACE, stroke=LIGHT, stroke_pt=1.2)
    case3_schema(s, 1, top=schema_y + 0.2)

    v_y = schema_y + schema_h + 0.15
    ocean_box(s, 0.55, v_y, 12.23, 0.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, 0.85, v_y, 11.6, 0.95,
             text="Весь транскрипт помещается в один запрос — усложнять сейчас означало бы "
                  "решать проблему, которой не существует",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Протоколы встреч", size=28)
    quote_block(s, 0.55, 1.15, 12.23, 1.2,
                "«На планёрке спросили: что мы решали по подрядчику X в мае?»",
                size=15.5)

    mid_y = 2.75
    mid_h = 1.55
    ocean_box(s, 0.55, mid_y, 12.23, mid_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.5)
    icon(s, "folder-search", "065A82", 96, 0.95, mid_y + (mid_h - 0.9) / 2, 0.9)
    text_box(s, 2.15, mid_y, 10.2, mid_h,
             text="Прошёл месяц — накопились десятки протоколов с разных встреч, по разным проектам",
             size=16, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    q_y = mid_y + mid_h + 0.3
    ocean_box(s, 0.55, q_y, 12.23, 1.15, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, 0.85, q_y, 11.6, 1.15, text="Меняет ли это архитектуру?", size=24,
             bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    """Case 3 verdict 2 = schema v2 (RAG expansion). v3: case3_schema
    (stage=2) -- v1 chain muted, RAG chain colored above it, feeding into
    the same "вызов модели" block. NO ladder_row. Iter-2 fix: schema frame
    now uses the SAME (schema_y, schema_h, top) as s22/s26 so row1 sits at
    an identical absolute y across all three slides -- this is what makes
    "the same schema growing" read as one continuous mutation rather than
    three different layouts. Row 2's slot (empty at this stage) is simply
    blank space inside the frame, matching s22's row0+row2 blank slots."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Протоколы встреч — разбор 1", size=22, y=0.2, h=0.4)
    quote_block(s, 0.55, 0.62, 12.23, 0.95,
                "«Растущий и меняющийся архив — вот когда нужен RAG»", size=14)

    schema_y, schema_h = 1.68, 4.3
    ocean_box(s, 0.55, schema_y, 12.23, schema_h, fill=SURFACE, stroke=LIGHT, stroke_pt=1.2)
    case3_schema(s, 2, top=schema_y + 0.2)

    v_y = schema_y + schema_h + 0.15
    ocean_box(s, 0.55, v_y, 12.23, 0.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, 0.85, v_y, 11.6, 0.95,
             text="Прежний путь никуда не делся — просто дополнен хранилищем, индексацией "
                  "и поиском перед тем же вызовом модели",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
    speaker_notes(s, load_notes("s24"))


def build_s25(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Протоколы встреч", size=28)
    quote_block(s, 0.55, 1.15, 12.23, 1.2,
                "«Поручения из протоколов мы всё равно руками переносим в таск-трекер...»",
                size=15.5)

    mid_y = 2.65
    mid_h = 1.4
    ocean_box(s, 0.55, mid_y, 12.23, mid_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.5)
    icon(s, "list-checks", "065A82", 96, 0.95, mid_y + (mid_h - 0.7) / 2, 0.7)
    text_box(s, 2.0, mid_y, 10.3, mid_h,
             text="Половина поручений из майских протоколов до сих пор без исполнителя — "
                  "руки до переноса просто не доходят",
             size=14.5, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.28)

    q_y = mid_y + mid_h + 0.3
    ocean_box(s, 0.55, q_y, 12.23, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, 0.85, text="Меняет ли это архитектуру ещё раз?",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s25"))


def build_s26(p):
    """Case 3 verdict 3 = schema v3 (agent expansion). v3: case3_schema
    (stage=3) -- v1+v2 muted, agent chain colored below, fed from the
    "протокол" block. NO ladder_row. Iter-2 fix: quote card was previously
    far too small (0.75in) for its text, causing horizontal overflow off
    the right edge of the slide -- resized to fit; schema now correctly
    fits its full 3-row footprint inside the available vertical band
    (title+quote reduced to a tighter footprint so all 3 rows + takeaway
    fit within 7.5in without the row2/row1 collision seen in iter 1)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Протоколы встреч — разбор 2", size=22, y=0.2, h=0.4)
    quote_block(s, 0.55, 0.62, 12.23, 0.95,
                "«Действие во внешней системе — нужен агент поверх RAG»", size=13.5)

    schema_y, schema_h = 1.68, 4.3
    ocean_box(s, 0.55, schema_y, 12.23, schema_h, fill=SURFACE, stroke=LIGHT, stroke_pt=1.2)
    case3_schema(s, 3, top=schema_y + 0.2)

    v_y = schema_y + schema_h + 0.15
    ocean_box(s, 0.55, v_y, 12.23, 0.95, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, 0.85, v_y, 11.6, 0.95,
             text="Начинайте с простейшего варианта, усложняйте только когда требование "
                  "не закрыть иначе",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """Quickfire setup with icon-scene. Iter-2 fix: enlarged scene + question
    boxes to close dead space at bottom."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Дайджест по чатам и почте", size=27)
    quote_block(s, 0.55, 1.15, 12.23, 1.05,
                "«Нужен еженедельный дайджест по рабочим чатам и почте команды»", size=15.5)

    scene_y = 2.55
    scene_h = 2.85
    ocean_box(s, 0.55, scene_y, 12.23, scene_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "message-square-quote", "065A82", 96, 0.95, scene_y + 0.7, 0.9)
    icon(s, "mail", "028090", 96, 2.35, scene_y + 0.7, 0.9)
    icon(s, "arrow-right", "F0AB00", 64, 3.6, scene_y + 0.9, 0.55)
    icon(s, "file-text", "065A82", 96, 4.5, scene_y + 0.7, 0.9)
    multipara_box(s, 5.85, scene_y + (scene_h - 1.6) / 2, 6.1, 1.6, [
        {"text": "Что обсуждали, что решили за неделю", "size": 17.5, "bold": True,
         "color": DEEP, "line_spacing": 1.25, "space_after": 10},
        {"text": "Дайджест нужен по пятницам, к 17:00", "size": 14.5, "color": SLATE,
         "line_spacing": 1.28},
    ])

    q_y = scene_y + scene_h + 0.3
    q_h = 7.15 - q_y
    ocean_box(s, 0.55, q_y, 12.23, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, q_h, text="Какая архитектура здесь нужна?",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Дайджест — разбор", size=28)

    grid_y = 1.5
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 1.95
    positive_card(s, 0.55, grid_y, cw, ch, "Агент — нужен", [
        {"text": "Обход нескольких источников по расписанию", "size": 13.5, "color": DEEP,
         "line_spacing": 1.25, "space_after": 6},
        {"text": "Многошаговая задача с обращением к внешним системам", "size": 13.5,
         "color": DEEP, "line_spacing": 1.25},
    ], icon_name="circle-check")

    negative_card(s, 0.55 + cw + gap, grid_y, cw, ch, "RAG — не нужен", [
        {"text": "Источники обходятся напрямую по расписанию, не по произвольному запросу",
         "size": 13.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Суммаризация — один вызов модели после сбора", "size": 13.5, "color": DEEP,
         "line_spacing": 1.25},
    ], icon_name="x-circle")

    footer_y = grid_y + ch + 0.35
    ocean_box(s, 0.55, footer_y, 12.23, 1.15, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, footer_y, 11.6, 1.15, text="Много источников — не то же самое, что RAG",
             size=20, bold=True, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """v3: REPLACED content -- Morgan Stanley (positive) + Replit (negative),
    Octomind demoted to a small footer note (per s29-rag-agent-in-prod.md)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "RAG и агент в проде", size=28)

    grid_y = 1.35
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 3.55
    positive_card(s, 0.55, grid_y, cw, ch, "AI @ Morgan Stanley Assistant", [
        {"text": "Партнёрство с OpenAI (март 2023), запуск в сентябре 2023",
         "size": 12.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "RAG по базе ~100 000 research-документов", "size": 12.5, "color": DEEP,
         "line_spacing": 1.25, "space_after": 6},
        {"text": "К середине 2024 — используют более 98% команд финансовых советников",
         "size": 12.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Доступность нужного документа: 20% → 80%", "size": 13, "bold": True,
         "color": DEEP, "line_spacing": 1.25},
    ], icon_name="circle-check")

    rx = 0.55 + cw + gap
    negative_card(s, rx, grid_y, cw, ch, "Replit, июль 2025", [
        {"text": "Агент удалил продакшен-базу (~1200 руководителей компаний затронуты) "
                 "во время code freeze, вопреки прямой инструкции",
         "size": 12.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Затем неверно утверждал, что откат невозможен — данные восстановили",
         "size": 12.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "После инцидента: разделение dev/prod + режим агента «только планирование»",
         "size": 12.5, "bold": True, "color": DEEP, "line_spacing": 1.25},
    ], icon_name="x-circle")

    footer_y = grid_y + ch + 0.22
    footer_note(s, "Ещё: Octomind — сняли LangChain после года в проде, вернулись к прямым "
                    "вызовам (17.06.2024)", y=footer_y)
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    build_divider(p, "s30", "4", "Внешний API или локальный инференс?", "Кейс: звонки продаж",
                  photo_path=SHOTS / "s30-divider-datacenter-real.jpg",
                  photo_credit="NASA · Wikimedia Commons · общественное достояние",
                  teaser="Это про то, где физически обрабатываются данные, а не какая модель лучше")


def build_s31(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Звонки продаж", size=28)
    quote_block(s, 0.55, 1.1, 12.23, 1.3,
                "«Хочу карточку по каждому звонку — потребность, возражения, договорённости»",
                size=15.5, role_icon="briefcase")

    photo_x, photo_y, photo_w, photo_h = 0.55, 2.6, 4.6, 3.95
    ocean_box(s, photo_x, photo_y, photo_w, photo_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    img_path = SHOTS / "s29-office-phonecall-real.jpg"
    if img_path.exists():
        pad = 0.12
        add_image_coverfit(s, img_path, photo_x + pad, photo_y + pad,
                            photo_w - 2 * pad, photo_h - 2 * pad)

    rx = photo_x + photo_w + 0.3
    rw = 12.23 - photo_w - 0.3
    ocean_box(s, rx, photo_y, rw, 1.85, fill=SURFACE, stroke=LIGHT)
    text_box(s, rx + 0.2, photo_y + 0.16, rw - 0.4, 1.55,
             text="8 менеджеров, по 15-20 звонков в день, записи лежат в телефонии",
             size=15, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    q_y = photo_y + 2.15
    q_h = photo_h - 2.15
    ocean_box(s, rx, q_y, rw, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.5)
    text_box(s, rx + 0.2, q_y, rw - 0.4, q_h, text="Как это построить?",
             size=20, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    speaker_notes(s, load_notes("s31",
                  extra="(Фото: OddibeKerfeld · Wikimedia Commons · CC BY-SA 3.0)"))


def build_s32(p):
    """v3: MERGED single slide -- lawyer's quote card at top + THREE
    jurisdiction cards (РФ/ЕС/США) below + one-line verdict at bottom (per
    s32-legal-map.md exact numbers). Was a single 420-ФЗ-only card in v2."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Звонки продаж", size=25, y=0.3, h=0.5)
    quote_block(s, 0.55, 0.85, 12.23, 1.05,
                "«Стоп. Вы эти записи в какой сервис отправляете?»", size=14.5,
                role_icon="scale")

    cards_y = 2.05
    ch = 3.15
    gap = 0.25
    n = 3
    cw = (12.23 - gap * (n - 1)) / n
    jurisdictions = [
        ("landmark", "РФ", [
            "152-ФЗ + 420-ФЗ (с 30.05.2025)",
            "Первая утечка: штраф 3-5 млн ₽",
            "Повторная: 1-3% годовой выручки (мин. 20 млн, макс. 500 млн ₽)",
        ]),
        ("scale", "ЕС", [
            "GDPR",
            "До €20 млн или 4% мирового оборота (что больше)",
            "Рекорд: Meta, €1,2 млрд, 2023",
        ]),
        ("building-2", "США", [
            "Единого федерального закона нет",
            "Штатные (Калифорния — CCPA/CPRA)",
            "Секторальные (например, HIPAA)",
        ]),
    ]
    for i, (ic, title, lines) in enumerate(jurisdictions):
        cx = 0.55 + i * (cw + gap)
        ocean_box(s, cx, cards_y, cw, ch, stroke=MID)
        icon(s, ic, "065A82", 96, cx + (cw - 0.5) / 2, cards_y + 0.18, 0.5)
        text_box(s, cx + 0.12, cards_y + 0.78, cw - 0.24, 0.4, text=title, size=17,
                 bold=True, color=MID, align=PP_ALIGN.CENTER)
        paras = [{"text": ln, "size": 11.5, "color": DEEP, "line_spacing": 1.25,
                  "space_after": 6} for ln in lines]
        multipara_box(s, cx + 0.18, cards_y + 1.3, cw - 0.36, ch - 1.4, paras)

    v_y = cards_y + ch + 0.2
    ocean_box(s, 0.55, v_y, 12.23, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, v_y, 11.6, 0.85,
             text="За сохранность персональных данных всерьёз борются все крупные юрисдикции",
             size=16, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s32"))


def build_s33(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Звонки продаж", size=28)

    grid_y = 1.5
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 2.5
    quote_block(s, 0.55, grid_y, cw, ch,
                "«GPU-сервера у нас нет, и бюджета на него в этом году нет»",
                size=16, role_icon="briefcase")
    text_box(s, 0.55 + 0.24, grid_y + ch - 0.4, cw - 0.48, 0.3, text="— CTO",
             size=12, italic=True, color=SLATE)

    rx = 0.55 + cw + gap
    quote_block(s, rx, grid_y, cw, ch,
                "«Мне не беседа нужна — пять полей в CRM по каждому звонку»",
                size=16, role_icon="user-round")
    text_box(s, rx + 0.24, grid_y + ch - 0.4, cw - 0.48, 0.3,
             text="— руководитель отдела продаж", size=12, italic=True, color=SLATE)

    q_y = grid_y + ch + 0.3
    ocean_box(s, 0.55, q_y, 12.23, 0.9, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, 0.9, text="Меняет ли это архитектуру ещё раз?",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s33"))


def build_s34(p):
    """v3: EXPANDED local-model detail -- T-lite/T-pro, 4-bit quantization,
    LoRA; open model list Qwen3/Llama/Gemma/Mistral/T-lite/T-pro. NO
    calls_ladder_row (v3 drops all ladder widgets)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Звонки продаж — разбор", size=26)

    v_y = 1.4
    v_h = 1.3
    ocean_box(s, 0.55, v_y, 12.23, v_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "cpu", "21295C", 64, 0.85, v_y + (v_h - 0.6) / 2, 0.6)
    text_box(s, 1.65, v_y + 0.12, 10.8, v_h - 0.24,
             text="Задача сузилась до шаблона — локальная модель 7-8 млрд параметров "
                  "справляется, данные не покидают периметр компании",
             size=14.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.22)

    models_y = v_y + v_h + 0.22
    models_h = 0.85
    ocean_box(s, 0.55, models_y, 12.23, models_h, fill=SURFACE, stroke=LIGHT)
    text_box(s, 0.75, models_y, 3.0, models_h,
             text="Открытые модели 7-8B:", size=12.5, bold=True, color=MID,
             anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 3.75, models_y, 8.7, models_h,
             text="Qwen3 · Llama · Gemma · Mistral · T-lite · T-pro (Т-Банк, Apache 2.0, на базе Qwen)",
             size=12.5, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    tech_y = models_y + models_h + 0.2
    tech_h = 0.85
    ocean_box(s, 0.55, tech_y, 12.23, tech_h, fill=SURFACE, stroke=LIGHT)
    text_box(s, 0.75, tech_y, 11.8, tech_h,
             text="4-bit квантование → 4-6 ГБ памяти, обычный сервер без GPU (Ollama / "
                  "llama.cpp) · тонкая подстройка — LoRA на одной GPU",
             size=12.5, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.22)

    lessons_y = tech_y + tech_h + 0.25
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 1.6
    ocean_box(s, 0.55, lessons_y, cw, ch, stroke=GOLD, stroke_pt=1.6, fill=GOLD_TINT)
    text_box(s, 0.55 + 0.2, lessons_y + 0.14, cw - 0.4, ch - 0.28,
             text="«Сложно для человека» ≠ «сложно для модели»",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    ocean_box(s, 0.55 + cw + gap, lessons_y, cw, ch, stroke=GOLD, stroke_pt=1.6, fill=GOLD_TINT)
    text_box(s, 0.55 + cw + gap + 0.2, lessons_y + 0.14, cw - 0.4, ch - 0.28,
             text="Чувствительность данных выясняется вопросами, а не предполагается",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    speaker_notes(s, load_notes("s34"))


def build_s35(p):
    """v3: REPLACED JetBrains -> Apple Intelligence + Google Gemini Nano."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Утечка в Samsung — и модели, которые уже в кармане", size=21)

    grid_y = 1.4
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 3.35
    negative_card(s, 0.55, grid_y, cw, ch, "Samsung, 2023", [
        {"text": "Разрешили ChatGPT 11.03.2023", "size": 12.5, "color": DEEP,
         "line_spacing": 1.25, "space_after": 6},
        {"text": "К 30.03 (~19 дней) — 3 утечки: 2× код производства, 1× транскрипт совещания",
         "size": 12.5, "bold": True, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Полный запрет генеративного ИИ на устройствах (май 2023)", "size": 12.5,
         "color": DEEP, "line_spacing": 1.25},
    ], icon_name="x-circle")

    rx = 0.55 + cw + gap
    positive_card(s, rx, grid_y, cw, ch, "Малые модели в кармане", [
        {"text": "Apple Intelligence — ~3 млрд параметров на устройстве (iOS 18.1, "
                 "октябрь 2024, iPhone 15 Pro+): суммаризация и правка текста офлайн",
         "size": 12, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Google Gemini Nano — 1,8-3,25 млрд параметров (Pixel/Android): "
                 "суммаризация в Recorder, Magic Compose — офлайн",
         "size": 12, "bold": True, "color": DEEP, "line_spacing": 1.25},
    ], icon_name="circle-check")

    takeaway_y = grid_y + ch + 0.25
    ocean_box(s, 0.55, takeaway_y, 12.23, 1.05, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, takeaway_y, 11.6, 1.05,
             text="Малые модели уже в кармане у половины аудитории",
             size=18, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s35"))


def build_s36(p):
    """Open invitation -- no worksheet, no numbered steps."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Ваши задачи", size=30, align=PP_ALIGN.LEFT)

    ill_y = 1.6
    ill_h = 3.5
    ocean_box(s, 0.55, ill_y, 12.23, ill_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "hand", "F0AB00", 96, 4.9, ill_y + 0.5, 1.2)
    icon(s, "users", "065A82", 96, 6.6, ill_y + 0.5, 1.2)
    text_box(s, 0.9, ill_y + 2.1, 11.5, 1.1,
             text="У кого есть похожая задача с работы — там, где стоял вопрос, нужен ли ИИ, и какой?",
             size=20, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.3)

    q_y = ill_y + ill_h + 0.3
    ocean_box(s, 0.55, q_y, 12.23, 0.9, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, 0.9,
             text="Достаточно сформулировать в двух-трёх предложениях, как её принёс бы заказчик",
             size=15.5, bold=True, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.2)
    speaker_notes(s, load_notes("s36"))


def build_s37(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Что унести", size=28)
    text_box(s, 0.55, 1.15, 12.23, 0.5,
             text="Семь уроков, по одному на кейс — ничего нового, только компиляция",
             size=15, italic=True, color=MID)

    lessons = [
        ("file-text", "Документы поставщиков", "Два примера — не выборка"),
        ("shield-alert", "Персоналка в логах", "Граница regex/модель — внутри задачи"),
        ("monitor", "Помощник поддержки", "Организационное ограничение решает не хуже техники"),
        ("store", "Описания товаров", "Встроенность бьёт отдельный продукт"),
        ("list-checks", "Протоколы встреч", "Архитектуру двигают требования, не мода"),
        ("mail", "Дайджест", "Много источников ≠ RAG"),
        ("phone", "Звонки продаж", "Сложно для человека ≠ сложно для модели"),
    ]
    grid_y = 1.85
    row_h = 0.65
    gap = 0.06
    for i, (ic, case, lesson) in enumerate(lessons):
        ry = grid_y + i * (row_h + gap)
        ocean_box(s, 0.55, ry, 12.23, row_h)
        icon(s, ic, "065A82", 64, 0.75, ry + (row_h - 0.34) / 2, 0.34)
        text_box(s, 1.3, ry, 3.5, row_h, text=case, size=13.5, bold=True, color=MID,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        filled_rect(s, 4.72, ry + row_h / 2 - 0.05, 0.1, 0.1, GOLD, radius=True,
                    radius_adj=0.5)
        text_box(s, 4.95, ry, 7.6, row_h, text=lesson, size=13.5, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    speaker_notes(s, load_notes("s37"))


def build_s38(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Домашнее чтение", size=28)
    text_box(s, 0.55, 1.15, 12.23, 0.45,
             text="Необязательно, но по каждой истории будет что обсудить",
             size=15, italic=True, color=MID)

    grid_y = 1.75
    gap = 0.28
    cw = (12.23 - gap * 2) / 3
    ch = 4.9
    cards = [
        ("handshake", "Сага Klarna, 2023-2025",
         "AI-ассистент поддержки: от впечатляющего старта до признания «мы зашли слишком "
         "далеко» и возврата к найму людей."),
        ("building-2", "NYC MyCity chatbot, 2024",
         "Городской чат-бот для малого бизнеса, дававший юридически некорректные советы. "
         "Расследование The Markup."),
        ("banknote", "Бонус, для настроения",
         "Автосалон подключил чат-бот на стороннем API без ограждений — и чат-бот "
         "«продал» внедорожник за один доллар."),
    ]
    for i, (ic, title, desc) in enumerate(cards):
        cx = 0.55 + i * (cw + gap)
        is_recommended = (i == 0)
        if is_recommended:
            ocean_box(s, cx, grid_y, cw, ch, stroke=GOLD, stroke_pt=2.0)
            chip(s, cx + cw - 1.35, grid_y + 0.2, 1.15, 0.34, "ГЛАВНОЕ", fill=GOLD,
                 color=DEEP, size=10.5)
        else:
            ocean_box(s, cx, grid_y, cw, ch)
        icon(s, ic, "065A82", 96, cx + (cw - 0.6) / 2, grid_y + 0.3, 0.6)
        text_box(s, cx + 0.2, grid_y + 1.15, cw - 0.4, 0.9, text=title, size=15, bold=True,
                 color=MID, align=PP_ALIGN.CENTER, line_spacing=1.2)
        text_box(s, cx + 0.2, grid_y + 2.15, cw - 0.4, ch - 2.35, text=desc, size=12.5,
                 color=DEEP, line_spacing=1.32)
    footer_note(s, "Если сегодня останется время — часть этого разберём вживую",
                y=6.85)
    speaker_notes(s, load_notes("s38"))


def build_s39(p):
    """Reserve slide -- no visible 'reserve' marker (frontmatter-only field)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Klarna, 2023-2025: одна история — четыре развилки", size=22)

    rows = [
        ("1", "ИИ или нет", "Типовые тикеты — формализуемы, хороший кандидат; но "
         "оптимизация по стоимости уронила качество"),
        ("2", "Встроить или своё", "Ассистент в существующем канале поддержки, внешний "
         "API (OpenAI) — не отдельный продукт"),
        ("3", "Архитектура", "Чат-ассистент с эскалацией к человеку, не автономный агент"),
        ("4", "Инференс", "Внешний облачный API с клиентскими данными поддержки"),
    ]
    ry = 1.3
    rh = 0.92
    gap = 0.1
    for num, label, desc in rows:
        ocean_box(s, 0.55, ry, 12.23, rh, stroke=TEAL)
        chip(s, 0.75, ry + (rh - 0.4) / 2, 0.4, 0.4, num, fill=MID, size=14)
        text_box(s, 1.3, ry + 0.1, 2.4, rh - 0.2, text=label, size=13.5, bold=True,
                 color=MID, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        text_box(s, 3.85, ry + 0.1, 8.75, rh - 0.2, text=desc, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
        ry += rh + gap

    facts_y = ry + 0.05
    ocean_box(s, 0.55, facts_y, 12.23, 0.85, fill=SURFACE, stroke=LIGHT)
    text_box(s, 0.8, facts_y, 11.7, 0.85,
             text="Запуск февраль 2024 · 30 дней: 2,3 млн чатов, 67% обращений автоматизировано · "
                  "заявлено «эквивалент 700 агентов»/$40М (2024) → «853 агента»/$60М (Q3 2025) — "
                  "знаменатель не раскрыт",
             size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    quote_y = facts_y + 1.0
    ocean_box(s, 0.55, quote_y, 12.23, 0.85, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, quote_y, 11.6, 0.85,
             text="«We went too far» — CEO Себастьян Семятковски, Bloomberg, май 2025",
             size=15, bold=True, italic=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s39"))


def build_s40(p):
    """Reserve quickfire setup with icon-scene. Iter-2 fix: enlarged scene +
    question boxes to close dead space at bottom (consistent with s07/s14/
    s27 fixes)."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Дедупликация новостей", size=28)
    quote_block(s, 0.55, 1.15, 12.23, 1.05,
                "«Поток новостей из многих источников — нужно склеивать одинаковые "
                "в одну карточку»", size=15)

    scene_y = 2.55
    scene_h = 2.85
    ocean_box(s, 0.55, scene_y, 12.23, scene_h, fill=SURFACE, stroke=TEAL, stroke_pt=1.6)
    icon(s, "message-square-quote", "065A82", 96, 0.95, scene_y + 0.65, 0.9)
    icon(s, "message-square-quote", "1C7293", 72, 2.15, scene_y + 0.95, 0.75)
    icon(s, "arrow-right", "F0AB00", 64, 3.25, scene_y + 1.05, 0.55)
    icon(s, "layers", "065A82", 96, 4.15, scene_y + 0.65, 0.9)
    multipara_box(s, 5.5, scene_y + (scene_h - 1.55) / 2, 6.4, 1.55, [
        {"text": "Десять почти одинаковых заголовков от разных изданий", "size": 17,
         "bold": True, "color": DEEP, "line_spacing": 1.28, "space_after": 10},
        {"text": "→ одна объединённая карточка события", "size": 14, "color": SLATE,
         "line_spacing": 1.28},
    ])

    q_y = scene_y + scene_h + 0.3
    q_h = 7.15 - q_y
    ocean_box(s, 0.55, q_y, 12.23, q_h, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, q_y, 11.6, q_h, text="Какая технология нужна?",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s40"))


def build_s41(p):
    s = blank(p)
    set_slide_bg(s, WHITE)
    slide_title(s, "Дедупликация — разбор", size=28)

    grid_y = 2.35
    gap = 0.3
    cw = (12.23 - gap) / 2
    ch = 1.95
    positive_card(s, 0.55, grid_y, cw, ch, "Fuzzy-matching / shingling — ~90%", [
        {"text": "Сравнение текстов по перекрывающимся фрагментам без обращения к модели",
         "size": 13.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Дёшево и предсказуемо", "size": 13.5, "bold": True, "color": DEEP,
         "line_spacing": 1.25},
    ], icon_name="filter")

    negative_card(s, 0.55 + cw + gap, grid_y, cw, ch, "LLM — только остаток", [
        {"text": "Два сообщения об одном событии написаны разными словами почти без общих "
         "фрагментов", "size": 13.5, "color": DEEP, "line_spacing": 1.25, "space_after": 6},
        {"text": "Перефраз без лексического пересечения", "size": 13.5, "italic": True,
         "color": DEEP, "line_spacing": 1.25},
    ], icon_name="sparkles")

    footer_y = grid_y + ch + 0.35
    ocean_box(s, 0.55, footer_y, 12.23, 1.15, fill=GOLD_TINT, stroke=GOLD, stroke_pt=1.4)
    text_box(s, 0.85, footer_y, 11.6, 1.15,
             text="Сначала дешёвый детерминированный фильтр — ИИ на то, что он не взял",
             size=18, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s41"))


def build_s42(p):
    """Hero closing -- full-bleed real photo. Kept from prior sessions (same
    NVIDIA GPU Wikimedia photo, same layout) after re-verification against
    s42-hero-closing.md wording; attempted a fresh 6-tier search for a
    literal "после работы" workspace-mood photo this session (see
    iteration-log.md) -- no better free-licensed candidate found after a
    genuine attempt in the time available, kept the well-documented GPU
    photo (bridges honestly to Lecture 2's "внутри модели" theme, which is
    hardware-adjacent even if not literally "рабочее пространство после
    работы")."""
    s = blank(p)
    set_slide_bg(s, DEEP)
    img_path = SHOTS / "s-closing-gpu-real.jpg"
    if img_path.exists():
        add_image(s, img_path, 0, 0, w=SLIDE_W_IN, h=SLIDE_H_IN)
    overlay = filled_rect(s, 0, 4.9, SLIDE_W_IN, 2.6, DEEP)
    overlay.fill.fore_color.rgb = DEEP
    try:
        alpha = etree.SubElement(overlay.fill.fore_color._xFill.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"),
            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
        alpha.set("val", "82000")
    except Exception:
        pass
    multipara_box(s, 0.6, 5.15, 11.6, 1.9, [
        {"text": "Сегодня вы решали, что строить.", "size": 26, "bold": True, "color": WHITE,
         "line_spacing": 1.2, "space_after": 4},
        {"text": "Лекция 2 — как модель устроена внутри", "size": 26, "bold": True, "color": GOLD,
         "line_spacing": 1.2},
    ])
    speaker_notes(s, load_notes("s42",
                  extra="(Фото: Mickael Courtiade · Wikimedia Commons · CC BY 2.0)"))


# ============================================================
# Orchestrate
# ============================================================

BUILDERS = [
    ("s01", build_s01), ("s02", build_s02), ("s03", build_s03), ("s04", build_s04),
    ("s05", build_s05), ("s06", build_s06), ("s07", build_s07), ("s08", build_s08),
    ("s09", build_s09), ("s10", build_s10), ("s11", build_s11), ("s12", build_s12),
    ("s13", build_s13), ("s14", build_s14), ("s15", build_s15), ("s16", build_s16),
    ("s17", build_s17), ("s18", build_s18), ("s19", build_s19), ("s20", build_s20),
    ("s21", build_s21), ("s22", build_s22), ("s23", build_s23), ("s24", build_s24),
    ("s25", build_s25), ("s26", build_s26), ("s27", build_s27), ("s28", build_s28),
    ("s29", build_s29), ("s30", build_s30), ("s31", build_s31), ("s32", build_s32),
    ("s33", build_s33), ("s34", build_s34), ("s35", build_s35), ("s36", build_s36),
    ("s37", build_s37), ("s38", build_s38), ("s39", build_s39), ("s40", build_s40),
    ("s41", build_s41), ("s42", build_s42),
]


def main():
    p = setup_pres()
    for sid, fn in BUILDERS:
        try:
            fn(p)
        except Exception as e:
            print(f"ERROR building {sid}: {e}")
            raise
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"Saved {OUT} — {len(BUILDERS)} slides")


if __name__ == "__main__":
    main()
