"""
Full 33-slide build of Лекции 5 «AI в финансовом секторе и ритейле».

Source-of-truth: deck.yaml + deck-part2.yaml (split >600 строк, loader reads
both) + chapter v2 finalized (3 части, ~22650 слов) + slides/*.md (33 файла:
32 LOCKED s01–s32 + s04a divider Раздела 1, readable speaker notes 150–300).

Issue #100 · worktree /tmp/lec-05-wt (branch phase-1-plan)

Palette LOCKED: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide. Motif «Ocean rounded box»
(radius 12, surface #F4F7FA, stroke #1C7293 1.5pt) на каждом content-слайде.

Canvas: 13.333" × 7.5" (16:9, [#55-1] patch). roadmap-bar ТОЛЬКО на
divider'ах (s04a/s10/s15/s20/s25) + cover (s02) — Л2-урок #40.

Render-style эталон: build_lec04.py (та же палитра/motif/типографика/
divider-шаблон/плотность). Build: python3 build_lec05.py → lec-05.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image

# === Palette (LOCKED) ===
DEEP = RGBColor(0x21, 0x29, 0x5C)
MID = RGBColor(0x06, 0x5A, 0x82)
LIGHT = RGBColor(0x1C, 0x72, 0x93)
TEAL = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF0, 0xAB, 0x00)
SLATE = RGBColor(0x5B, 0x66, 0x78)
COVER_OUTLINE = RGBColor(0xD9, 0xE2, 0xEC)
GOLD_TINT = RGBColor(0xFD, 0xF3, 0xDC)
TEAL_TINT = RGBColor(0xE4, 0xF1, 0xF2)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)

# === Constants ===
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path("/tmp/lec-05-wt/library/lectures/lec-05")
ASSETS = ROOT / "rendered/assets"
ICONS = ASSETS / "icons"
CHARTS = ASSETS / "charts"
DIAGRAMS = ASSETS / "diagrams"
SLIDES_DIR = ROOT / "slides"
OUT = ROOT / "rendered/lec-05.pptx"
FONT_HEAD = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Liberation Mono"


# ============================================================
# Helpers (architecture mirrors lec-04 build — proven)
# ============================================================

def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def disable_shadow(shp):
    sppr = shp._element.spPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for el in sppr.findall(ns + "effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, ns + "effectLst")


def text_box(slide, x, y, w, h, text, *,
             size=16, bold=False, italic=False, color=DEEP,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    pgr = tf.paragraphs[0]
    pgr.alignment = align
    pgr.line_spacing = line_spacing
    r = pgr.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def text_runs(slide, x, y, w, h, runs, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15, font=FONT_BODY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    pgr = tf.paragraphs[0]
    pgr.alignment = align
    pgr.line_spacing = line_spacing
    for cfg in runs:
        if cfg.get("newpara"):
            pgr = tf.add_paragraph()
            pgr.alignment = cfg.get("align", align)
            pgr.line_spacing = cfg.get("line_spacing", line_spacing)
            if cfg.get("space_before") is not None:
                pgr.space_before = Pt(cfg["space_before"])
        r = pgr.add_run()
        r.text = cfg["text"]
        r.font.name = cfg.get("font", font)
        r.font.size = Pt(cfg.get("size", 16))
        r.font.bold = cfg.get("bold", False)
        r.font.italic = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", DEEP)
    return tb


def ocean_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_pt=1.5,
              radius_pt=12.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        adj = max(0.035, min(0.22,
                             (radius_pt / 72.0) / max(min(w, h) / 2.0, 0.5)))
        shp.adjustments[0] = adj
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def filled_rect(slide, x, y, w, h, fill, *, stroke=None, stroke_pt=0.0,
                radius=False, radius_adj=0.16):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    if radius:
        try:
            shp.adjustments[0] = radius_adj
        except Exception:
            pass
    disable_shadow(shp)
    return shp


def right_arrow(slide, x, y, w, h, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.0)
    disable_shadow(shp)
    return shp


def circle(slide, x, y, d, fill, *, stroke=None, stroke_pt=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_pt)
    disable_shadow(shp)
    return shp


def chip(slide, x, y, w, h, text, *, fill=MID, stroke=None, color=WHITE,
         size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    pgr = tf.paragraphs[0]
    pgr.alignment = PP_ALIGN.CENTER
    r = pgr.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    disable_shadow(shp)
    return shp


def connector(slide, x1, y1, x2, y2, color=LIGHT, width=2.0, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = etree.SubElement(
            ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                "prstDash")
        pd.set("val", dash)
    return cn


def add_image(slide, path, x, y, w=None, h=None, preserve_aspect=True):
    """[#73-render-1] aspect-safe: pass only the constraining dimension."""
    path = Path(path)
    if not path.exists():
        return
    if preserve_aspect and w is not None and h is not None:
        try:
            img = Image.open(path)
            iw, ih = img.size
            img.close()
        except Exception:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                     width=Inches(w))
            return
        ir = iw / ih
        br = w / h
        if ir > br:
            ah = w / ir
            slide.shapes.add_picture(str(path), Inches(x),
                                     Inches(y + (h - ah) / 2), width=Inches(w))
        else:
            aw = h * ir
            slide.shapes.add_picture(str(path), Inches(x + (w - aw) / 2),
                                     Inches(y), height=Inches(h))
    elif w is not None and h is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def slide_title(slide, text, *, y=0.40, h=0.92, w=12.25, x=0.55, size=26,
                color=DEEP, bold=True, line_spacing=1.10, align=PP_ALIGN.LEFT):
    text_box(slide, x=x, y=y, w=w, h=h, text=text,
             size=size, bold=bold, color=color, line_spacing=line_spacing,
             align=align)


def gold_callout(slide, x, y, w, h, text, *, size=15, bold=True,
                 color=DEEP, align=PP_ALIGN.LEFT, stroke_pt=1.5):
    filled_rect(slide, x, y, w, h, GOLD_TINT, stroke=GOLD, stroke_pt=stroke_pt,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.20)


def teal_callout(slide, x, y, w, h, text, *, size=14, bold=False,
                 color=DEEP, align=PP_ALIGN.LEFT):
    filled_rect(slide, x, y, w, h, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                radius=True, radius_adj=0.10)
    text_box(slide, x=x + 0.24, y=y + 0.06, w=w - 0.48, h=h - 0.12, text=text,
             size=size, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE,
             align=align, line_spacing=1.18)


def footer(slide, text):
    text_box(slide, x=0.55, y=7.04, w=12.25, h=0.34, text=text,
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.LEFT,
             line_spacing=1.0)


def stat_plate(slide, x, y, w, h, label, value, *, highlight=False,
               vsize=25, lsize=12.5):
    bg = GOLD_TINT if highlight else SURFACE
    edge = GOLD if highlight else SOFT_GREY
    filled_rect(slide, x, y, w, h, bg, stroke=edge,
                stroke_pt=(1.5 if highlight else 1.0),
                radius=True, radius_adj=0.10)
    text_box(slide, x + 0.22, y + 0.08, w - 0.44, 0.30, label,
             size=lsize, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, x + 0.22, y + 0.34, w - 0.44, h - 0.42, value,
             size=vsize, bold=True, color=(GOLD if highlight else DEEP),
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


def icon(slide, name, x, y, size, variant="mid"):
    add_image(slide, ICONS / f"{name}-{variant}.png", x, y, size, size)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def load_notes(slide_id):
    files = list(SLIDES_DIR.glob(f"{slide_id}-*.md"))
    if not files:
        return ""
    md = files[0].read_text(encoding="utf-8")
    m = re.search(r'## Speaker notes\s*\n(.*?)(?=\n## |\n---\s*\n## |\Z)',
                  md, re.DOTALL)
    notes = m.group(1).strip() if m else ""
    notes = re.sub(r'\n+---\s*$', '', notes)
    return notes.strip()


# ============================================================
# Deck loader — deck.yaml split на 2 части (≤600 строк каждая).
# Loader читает ОБЕ части, объединяет slides, вставляет s04a после s04.
# ============================================================
def load_deck():
    try:
        import yaml
    except ImportError:
        return None
    d1 = yaml.safe_load((ROOT / "deck.yaml").read_text(encoding="utf-8"))
    d2 = yaml.safe_load((ROOT / "deck-part2.yaml").read_text(encoding="utf-8"))
    slides = list(d1.get("slides", [])) + list(d2.get("slides", []))
    ids = [s["id"] for s in slides]
    expected = [f"s{n:02d}" for n in range(1, 33)]
    expected.insert(4, "s04a")  # s04a after s04
    assert ids == expected, (
        f"deck slide order mismatch:\n got={ids}\n exp={expected}")
    tot = d2.get("totals", {}).get("slides")
    assert tot == 33, f"deck-part2 totals.slides={tot}, expected 33"
    return {"slides": slides, "totals": d2.get("totals", {}),
            "deck": d1.get("deck", {})}


# ============================================================
# Section divider — unified template (7-card roadmap, gold current)
# Sections of Лекции 5: 0..6.
# ============================================================
NAV = [
    ("0", "Открытие"),
    ("1", "Прогноз"),
    ("2", "Аномалии"),
    ("3", "Скоринг"),
    ("4", "LLM"),
    ("5", "Рекоменд."),
    ("6", "Фреймворк"),
]


def roadmap_bar(slide, here_idx, *, y=6.50):
    """7-card progress bar; current section gold-bordered."""
    n = len(NAV)
    gap = 0.12
    bx = 0.55
    total_w = 12.25
    cw = (total_w - gap * (n - 1)) / n
    ch = 0.60
    for i, (num, label) in enumerate(NAV):
        x = bx + i * (cw + gap)
        cur = (i == here_idx)
        if cur:
            filled_rect(slide, x, y, cw, ch, GOLD_TINT, stroke=GOLD,
                        stroke_pt=2.0, radius=True, radius_adj=0.12)
        else:
            filled_rect(slide, x, y, cw, ch, SURFACE, stroke=SOFT_GREY,
                        stroke_pt=1.0, radius=True, radius_adj=0.12)
        text_box(slide, x=x + 0.04, y=y + 0.07, w=cw - 0.08, h=0.22,
                 text=f"Раздел {num}", size=10, bold=True,
                 color=(DEEP if cur else LIGHT), align=PP_ALIGN.CENTER)
        text_box(slide, x=x + 0.04, y=y + 0.30, w=cw - 0.08, h=0.26,
                 text=label, size=10.5, bold=cur,
                 color=(DEEP if cur else SLATE), align=PP_ALIGN.CENTER,
                 line_spacing=0.95)


def build_section_divider(p, here_idx, subtitle, bridge, sid):
    """Distinct divider (NO ocean motif): giant decorative section digit on
    the right, РАЗДЕЛ N + subtitle + 1-line narrative bridge on the left,
    gold-current roadmap bar at bottom."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=8.45, y=0.30, w=4.6, h=6.0, text=str(here_idx),
             size=400, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, x=0.75, y=1.55, w=7.4, h=0.55,
             text=f"РАЗДЕЛ {here_idx}", size=20, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.18, 0.70, 0.05, fill=GOLD)
    text_box(s, x=0.75, y=2.55, w=7.6, h=1.85, text=subtitle,
             size=36, bold=True, color=DEEP, line_spacing=1.08)
    text_box(s, x=0.78, y=4.62, w=7.5, h=1.55, text=bridge,
             size=18, italic=True, color=LIGHT, line_spacing=1.22)
    roadmap_bar(s, here_idx, y=6.50)
    speaker_notes(s, load_notes(sid))
    return s


# ============================================================
# Slide builders — 33 slides (32 LOCKED s01..s32 + s04a divider)
# ============================================================

def build_s01(p):
    """case_study — Zillow iBuying hook. Left: 3 numbers in ocean box.
    Right: frame «какой тип ИИ и почему не LLM». Gold: $500M+ закрыто."""
    s = blank(p)
    slide_title(s, "Одна прогнозная модель закрыла целое направление крупной "
                   "компании.", size=23, w=11.0, h=0.96)
    icon(s, "landmark", 12.05, 0.42, 0.74, "mid")
    # LEFT — mega-stat dominates; 3 small support stats below
    lx, ly, lw, lh = 0.55, 1.46, 7.05, 4.55
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.30, ly + 0.18, lw - 0.60, 0.30,
             "Zillow Offers · iBuying · ноябрь 2021",
             size=14, bold=True, color=MID)
    # mega number — the single point of impact
    filled_rect(s, lx + 0.30, ly + 0.54, lw - 0.60, 1.78, GOLD_TINT,
                stroke=GOLD, stroke_pt=2.0, radius=True, radius_adj=0.08)
    text_box(s, lx + 0.50, ly + 0.66, lw - 1.00, 1.02, "$500M+",
             size=68, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.50, ly + 1.70, lw - 1.00, 0.52,
             "совокупные потери — целое направление закрыто",
             size=15, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    sup = [("$300M+", "списание за квартал"),
           ("~2000", "уволено (≈25%)"),
           ("−25%", "акция за дни")]
    n = 3
    gap = 0.18
    cw = (lw - 0.60 - gap * (n - 1)) / n
    sx = lx + 0.30
    for val, lab in sup:
        filled_rect(s, sx, ly + 2.50, cw, 0.94, SURFACE, stroke=SOFT_GREY,
                    stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, sx + 0.06, ly + 2.58, cw - 0.12, 0.42, val,
                 size=21, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, sx + 0.06, ly + 3.00, cw - 0.12, 0.38, lab,
                 size=11, color=MID, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        sx += cw + gap
    text_box(s, lx + 0.30, ly + 3.62, lw - 0.60, 0.78,
             "Модель прогнозировала цену дома · Zillow авто покупала, "
             "ремонтировала, перепродавала тысячи домов · систематически "
             "переоценивала",
             size=11.5, color=SLATE, line_spacing=1.16)
    # RIGHT — "не ChatGPT" + open question
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 2.30, fill=TEAL_TINT, stroke=TEAL, stroke_pt=2.0)
    icon(s, "circle-help", rx + 0.28, ly + 0.24, 0.50, "teal")
    text_box(s, rx + 0.92, ly + 0.28, rw - 1.10, 0.42, "Это был не ChatGPT",
             size=17, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, rx + 0.28, ly + 0.90, rw - 0.56, 1.28,
             "Прогнозная модель оценивала число (цену) по табличным и "
             "геоданным. LLM здесь не применялась и не могла — задача не "
             "текстовая.",
             size=13.5, color=DEEP, line_spacing=1.20)
    ocean_box(s, rx, ly + 2.46, rw, 2.09)
    text_box(s, rx + 0.26, ly + 2.62, rw - 0.52, 1.80,
             "Какой тип ИИ это был — и почему обычная ошибка модели стала "
             "потерей бизнеса, а не мелкой неточностью?",
             size=14.5, italic=True, color=MID, line_spacing=1.22,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 6.20, 12.25, 0.66,
                 "Ошибка прогноза нормальна. Бизнес рухнул из-за того, КУДА "
                 "она подключена: к авто-, необратимому действию в масштабе "
                 "тысяч домов.", size=14)
    speaker_notes(s, load_notes("s01"))


def build_s02(p):
    """cover — distinct, NO ocean motif. Mega «05» + title + roadmap-bar."""
    s = blank(p)
    set_slide_bg(s, SURFACE)
    text_box(s, x=7.7, y=0.95, w=5.7, h=5.2, text="05",
             size=300, bold=True, color=COVER_OUTLINE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text_box(s, x=0.75, y=1.45, w=6.6, h=0.5, text="ЛЕКЦИЯ 5",
             size=18, bold=True, color=TEAL)
    filled_rect(s, 0.78, 2.02, 0.70, 0.05, fill=TEAL)
    text_box(s, x=0.75, y=2.42, w=7.7, h=2.5,
             text="AI в финансовом\nсекторе и ритейле",
             size=44, bold=True, color=DEEP, line_spacing=1.10)
    filled_rect(s, 0.78, 5.18, 0.05, 0.56, fill=TEAL)
    text_box(s, x=1.02, y=5.16, w=7.6, h=0.62,
             text="Под какую задачу — какой тип ИИ,\nпочему именно он, и где "
                  "он ломается",
             size=18, color=MID, line_spacing=1.18)
    roadmap_bar(s, 0, y=6.50)
    speaker_notes(s, load_notes("s02"))


def build_s03(p):
    """comparison — bridge Л2/Л4 → палитра 5 типов. KEYSTONE."""
    s = blank(p)
    slide_title(s, "Один тип ИИ вглубь (Л4) → палитра типов под палитру "
                   "задач (Л5).", size=24, y=0.32, h=0.56, w=12.25)
    # compact secondary bridge strip (Л2/Л4 → Л5) — single line each
    bx, by, bw, bh = 0.55, 0.98, 12.25, 1.16
    ocean_box(s, bx, by, bw, bh)
    bridges = [
        ("Л2: дерево «какой ИИ» — ML / глубокое / LLM",
         "Л5: пять структурно разных типов под разные задачи"),
        ("Л4: один тип (LLM-кодер) в глубину по автономности",
         "Л5: палитра типов под палитру задач одной индустрии"),
    ]
    yy = by + 0.16
    for a, b in bridges:
        text_box(s, bx + 0.26, yy, 5.55, 0.40, a,
                 size=11.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        text_box(s, bx + 5.86, yy + 0.06, 0.34, 0.30, "→",
                 size=15, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, bx + 6.26, yy, bw - 6.52, 0.40, b,
                 size=11.5, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        yy += 0.44
    # MAIN VISUAL — palette of 5 types, large
    text_box(s, 0.55, 2.34, 12.25, 0.34,
             "Карта пяти типов ИИ под пять разных задач (+ сквозной "
             "CV-пласт):", size=15, bold=True, color=MID)
    types = ["Прогноз\nрядов", "Поиск\nаномалий", "ML-\nскоринг",
             "LLM-\nассистенты", "Рекомендации\nи цены"]
    icns = ["trending-up", "radar", "scale", "message-circle",
            "shopping-cart"]
    n = 5
    gap = 0.22
    cw = (12.25 - gap * (n - 1)) / n
    ax, ay, ah = 0.55, 2.78, 2.06
    for i in range(n):
        ocean_box(s, ax, ay, cw, ah)
        icon(s, icns[i], ax + cw / 2 - 0.36, ay + 0.30, 0.72, "mid")
        for li, line in enumerate(types[i].split("\n")):
            text_box(s, ax + 0.04, ay + 1.18 + li * 0.36, cw - 0.08, 0.36,
                     line, size=14, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.0)
        ax += cw + gap
    gold_callout(s, 0.55, 5.10, 12.25, 1.66,
                 "Большую часть ценности здесь даёт НЕ языковая модель. "
                 "LLM — не универсальный молоток: разная задача требует "
                 "разного типа ИИ. Это и есть «палитра типов под палитру "
                 "задач».", size=16)
    speaker_notes(s, load_notes("s03"))


def build_s04(p):
    """assertion_visual — central question + 5-step pattern + disclaimer."""
    s = blank(p)
    slide_title(s, "Центральный вопрос лекции.",
                size=27, y=0.34, h=0.58, w=9.5)
    icon(s, "target", 12.05, 0.36, 0.78, "gold")
    bx, by, bw, bh = 0.55, 1.18, 12.25, 1.62
    ocean_box(s, bx, by, bw, bh, fill=GOLD_TINT, stroke=GOLD, stroke_pt=2.0)
    text_box(s, bx + 0.40, by + 0.18, bw - 0.80, bh - 0.36,
             "«Финансы и ритейл — отрасли максимального внедрения ИИ. Под "
             "какую задачу — какой тип ИИ, почему именно он (а не LLM "
             "везде), и где этот тип ломается?»",
             size=20, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16, align=PP_ALIGN.CENTER)
    # 5-step pattern ribbon
    text_box(s, 0.55, 2.96, 12.25, 0.30,
             "Каждый тип разбираем по одной схеме (5 шагов):",
             size=14, bold=True, color=MID)
    steps = ["задача", "тип ИИ\nи почему он", "реальный\nпример",
             "где\nломается", "альтернатива\nи критерий"]
    n = 5
    gap = 0.18
    cw = (12.25 - gap * (n - 1)) / n
    sx, sy, sh = 0.55, 3.30, 0.92
    for i in range(n):
        filled_rect(s, sx, sy, cw, sh, TEAL_TINT, stroke=TEAL, stroke_pt=1.5,
                    radius=True, radius_adj=0.12)
        text_box(s, sx + 0.06, sy, cw - 0.12, sh, f"{i + 1}. {steps[i]}",
                 size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
        if i < n - 1:
            text_box(s, sx + cw - 0.02, sy, 0.20, sh, "→",
                     size=16, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        sx += cw + gap
    teal_callout(s, 0.55, 4.42, 12.25, 0.66,
                 "Карта — это структура лекции, не требование понять всё "
                 "сейчас. Прогноз · аномалии · скоринг · LLM · "
                 "рекомендации (+ сквозной CV).",
                 size=13.5, bold=True)
    gold_callout(s, 0.55, 5.26, 12.25, 1.50,
                 "Ответ — не «ИИ хорош» и не «ИИ плох», а аппарат: назвать "
                 "уместный тип, обосновать выбор и заранее увидеть точку "
                 "обязательного человека. К концу схема станет вашим "
                 "инструментом выбора.", size=14)
    speaker_notes(s, load_notes("s04"))


def build_s04a(p):
    build_section_divider(
        p, 1, "Прогнозирование:\nспрос, продажи, отток",
        "Начнём с самой массовой задачи ритейла — и сразу с типа ИИ, "
        "который НЕ LLM.", "s04a")


def build_s05(p):
    """assertion_visual — prognosis vs LLM, 2-col parallel."""
    s = blank(p)
    slide_title(s, "Прогноз ряда и языковая модель решают структурно разные "
                   "задачи.", size=24, y=0.34, h=0.62, w=12.25)
    bx, by, bw, bh = 0.55, 1.20, 12.25, 3.40
    ocean_box(s, bx, by, bw, bh)
    col_w = [6.0, 6.05]
    heads = ["Языковая модель (LLM)", "Прогноз временного ряда"]
    cx = bx + 0.20
    for j, hd in enumerate(heads):
        filled_rect(s, cx, by + 0.18, col_w[j] - 0.20, 0.48,
                    MID if j == 0 else TEAL)
        text_box(s, cx, by + 0.18, col_w[j] - 0.20, 0.48, hd,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
    rows = [
        ("предсказывает следующий ТОКЕН текста",
         "предсказывает следующее ЧИСЛО в ряду"),
        ("оптимизирована на правдоподобие формулировки",
         "нужна числовая точность + калиброванная неопределённость"),
        ("нет внутреннего понятия «сезонность» / «тренд»",
         "ряд = тренд + сезонность + шум"),
        ("видела тексты О продажах",
         "видит сам ряд продаж как ряд"),
    ]
    yy = by + 0.78
    for a, b in rows:
        cx = bx + 0.20
        for j, cc in enumerate((a, b)):
            filled_rect(s, cx, yy, col_w[j] - 0.20, 0.58,
                        SURFACE, stroke=SOFT_GREY, stroke_pt=0.75)
            text_box(s, cx + 0.16, yy, col_w[j] - 0.42, 0.58, cc,
                     size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.05)
            cx += col_w[j]
        yy += 0.62
    teal_callout(s, 0.55, 4.74, 12.25, 0.62,
                 "Тип ИИ под прогноз: классическая статистика и табличный "
                 "ML — ARIMA-семейство, градиентный бустинг. НЕ "
                 "генеративная модель, НЕ LLM.",
                 size=13.5, bold=True)
    gold_callout(s, 0.55, 5.52, 12.25, 1.26,
                 "Поставить «спрогнозируй спрос» текстовой модели — она "
                 "выдаст правдоподобное число без связи с сезонностью: "
                 "галлюцинация в числовой обёртке. Тип ИИ определяется "
                 "СТРУКТУРОЙ задачи, а не модностью инструмента.", size=14)
    speaker_notes(s, load_notes("s05"))


def build_s06(p):
    """assertion_visual — 3 argument cards (why not LLM) + alternative."""
    s = blank(p)
    slide_title(s, "Почему здесь табличная модель, а не LLM — три "
                   "аргумента, а не оговорка.", size=22, y=0.34, h=0.58,
                w=12.25)
    cards = [
        ("database", "1. Структура данных",
         "табличные числовые ряды (дата, точка, товар, цена, признаки). "
         "Модели для рядов извлекают тренд и сезонность; перевод ряда в "
         "текст для LLM теряет именно эту структуру.", "mid", False),
        ("activity", "2. Измеримость и калибровка",
         "решение о размере закупки строится на численной метрике ошибки "
         "+ доверительном интервале. Классические модели дают это из "
         "коробки; от LLM такого нет.", "mid", False),
        ("triangle-alert", "3. Что сломается",
         "LLM даст правдоподобное число без обоснованной неопределённости "
         "— а решение всё равно придётся принять, и цена систематической "
         "ошибки × объём (механизм Zillow).", "gold", True),
    ]
    n = 3
    gap = 0.26
    cw = (12.25 - gap * (n - 1)) / n
    cx = 0.55
    cy, chh = 1.22, 2.74
    for ic, ttl, body, var, hi in cards:
        if hi:
            ocean_box(s, cx, cy, cw, chh, fill=GOLD_TINT, stroke=GOLD,
                      stroke_pt=2.0)
        else:
            ocean_box(s, cx, cy, cw, chh)
        icon(s, ic, cx + 0.24, cy + 0.22, 0.50, var)
        text_box(s, cx + 0.86, cy + 0.26, cw - 1.05, 0.50, ttl,
                 size=15, bold=True, color=(DEEP if hi else MID),
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, cx + 0.26, cy + 0.90, cw - 0.52, chh - 1.02, body,
                 size=12, color=DEEP, line_spacing=1.16)
        cx += cw + gap
    teal_callout(s, 0.55, 4.12, 12.25, 0.84,
                 "Альтернатива и критерий: правильный инструмент — "
                 "табличная прогнозная модель с измеримой ошибкой. Даже её "
                 "недостаточно при нестационарной среде + необратимом "
                 "капитальном действии без человеческого контроля.",
                 size=13.5, bold=True)
    gold_callout(s, 0.55, 5.08, 12.25, 0.86,
                 "Задача: прогноз спроса (полка / закупка), cash-flow "
                 "(ликвидность), оттока клиентов (кого удержать).", size=14)
    footer(s, "Если ряды содержат PII граждан РФ — отправлять их в публичный "
              "облачный LLM нарушает требования локализации (ФЗ-152); "
              "табличная модель разворачивается on-premise.")
    speaker_notes(s, load_notes("s06"))


def build_s07(p):
    """case_study — X5 + Магнит, 2 panels."""
    s = blank(p)
    slide_title(s, "Прогноз диктуется задачей: ни X5, ни Магнит не взяли "
                   "LLM.", size=24, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.30, 6.10, 3.95
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.34,
             "X5 Group · «Пятёрочка», «Перекрёсток»",
             size=14, bold=True, color=MID)
    x5 = [("Точность прогноза спроса", "> 70%", False),
          ("Доп. выручка от ML-инструментов", "+5 млрд ₽", True),
          ("Списания просрочки", "−2%", False)]
    cyy = ly + 0.58
    for lab, val, hi in x5:
        bg = GOLD_TINT if hi else SURFACE
        filled_rect(s, lx + 0.26, cyy, lw - 0.52, 0.96, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.44, cyy + 0.10, lw - 2.30, 0.76, lab,
                 size=12.5, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
        text_box(s, lx + lw - 2.05, cyy + 0.06, 1.66, 0.84, val,
                 size=22, bold=True, color=(GOLD if hi else DEEP),
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        cyy += 1.04
    rx, rw = 6.85, 5.95
    ocean_box(s, rx, ly, rw, 3.95)
    text_box(s, rx + 0.26, ly + 0.16, rw - 0.52, 0.34,
             "Магнит — импортозамещение прогнозных систем",
             size=14, bold=True, color=MID)
    mg = [
        ("До 2022:", "иностранные вендоры класса SAP / Blue Yonder"),
        ("После ухода:", "собственная система прогноза + автозаказа (F&R)"),
        ("Статус:", "пилот на распределительном центре (2024–2025)"),
    ]
    fy = ly + 0.62
    for a, b in mg:
        text_box(s, rx + 0.26, fy, 1.55, 0.66, a, size=12.5, bold=True,
                 color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 1.86, fy, rw - 2.12, 0.66, b, size=12.5, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        fy += 0.74
    teal_callout(s, rx, ly + 2.92, rw, 0.90,
                 "Уход иностранных вендоров сделал импортозамещение "
                 "прогнозных систем прямой инженерной задачей отрасли.",
                 size=12.5, bold=False)
    gold_callout(s, 0.55, 5.40, 12.25, 0.88,
                 "Ни X5, ни Магнит не решают задачу языковой моделью — они "
                 "строят специализированные прогнозные системы, потому что "
                 "ТИП ИИ ДИКТУЕТСЯ ЗАДАЧЕЙ.", size=14)
    footer(s, "По данным компаний и отраслевых обзоров (X5/TAdviser 2023; "
              "shoppers.media/TAdviser 2024–2025); числа — заявленные "
              "компаниями, сверяются в день лекции.")
    speaker_notes(s, load_notes("s07"))


def build_s08(p):
    """process — time-series decomposition diagram + товаровед analogy."""
    s = blank(p)
    slide_title(s, "Прогноз продлевает паттерны прошлого — здесь его сила и "
                   "уязвимость.", size=23, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.22, 7.30, 3.95
    ocean_box(s, lx, ly, lw, lh)
    add_image(s, DIAGRAMS / "d08-timeseries-decomp.png",
              lx + 0.20, ly + 0.16, lw - 0.40, lh - 0.62)
    text_box(s, lx + 0.24, ly + lh - 0.42, lw - 0.48, 0.34,
             "Любой ряд продаж = тренд + сезонность + шум; модель "
             "продлевает регулярную часть.",
             size=10.5, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    rx, rw = 8.05, 4.75
    ocean_box(s, rx, ly, rw, 3.95)
    text_box(s, rx + 0.24, ly + 0.16, rw - 0.48, 0.30,
             "Что делает модель — численно", size=13.5, bold=True,
             color=MID)
    pts = [
        "раскладывает историю на тренд + сезонность + шум",
        "обучается на регулярных паттернах, продлевает их",
        "добавляет поправки на промо / праздник",
    ]
    py = ly + 0.54
    for t in pts:
        circle(s, rx + 0.26, py + 0.05, 0.12, MID)
        text_box(s, rx + 0.52, py, rw - 0.78, 0.56, t,
                 size=12, color=DEEP, line_spacing=1.10)
        py += 0.58
    text_box(s, rx + 0.24, py + 0.06, rw - 0.48, 0.46,
             "Прогноз → действие: «спрос ≈ N» → «заказать N + запас»",
             size=12, bold=True, color=TEAL, line_spacing=1.08)
    text_box(s, rx + 0.24, py + 0.62, rw - 0.48, 1.00,
             "Аналогия: модель прогноза — как опытный товаровед «к выходным "
             "возьми больше», но для миллионов пар магазин × товар сразу.",
             size=11.5, italic=True, color=SLATE, line_spacing=1.16)
    gold_callout(s, 0.55, 5.32, 12.25, 0.96,
                 "Встроенная уязвимость: прогноз силён ровно настолько, "
                 "насколько ЗАВТРА ПОХОЖЕ НА ВЧЕРА. Меняется среда "
                 "качественно — модель уверенно экстраполирует паттерны, "
                 "которых больше нет. Ключ к следующему слайду.", size=14)
    speaker_notes(s, load_notes("s08"))


def build_s09(p):
    """case_study — Zillow payoff: distribution shift + asymmetry + Knight."""
    s = blank(p)
    slide_title(s, "Разорил Zillow не дрейф модели, а то, к чему был "
                   "подключён её выход.", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.16, 7.05, 4.10
    ocean_box(s, lx, ly, lw, lh)
    blocks = [
        ("distribution shift (дрейф распределения)",
         "данные в реальности перестали быть похожи на обучающие — рынок "
         "жилья 2020–2021 стал статистически другим"),
        ("асимметрия цены ошибки",
         "ошибка рекомендации ≈ 0 · та же ошибка, на которой куплен дом = "
         "десятки тыс. $ × N, и НЕОБРАТИМА"),
        ("Три вещи вместе = фатально",
         "необратимое действие × авто по выходу модели × нестационарная "
         "среда без circuit-breaker (авто-стоп при аномалии)"),
    ]
    by = ly + 0.20
    for ttl, body in blocks:
        text_box(s, lx + 0.26, by, lw - 0.52, 0.28, ttl,
                 size=13, bold=True, color=MID)
        text_box(s, lx + 0.26, by + 0.30, lw - 0.52, 0.60, body,
                 size=12, color=DEEP, line_spacing=1.12)
        by += 1.02
    filled_rect(s, lx + 0.26, by + 0.04, lw - 0.52, 0.62, TEAL_TINT,
                stroke=TEAL, stroke_pt=1.2, radius=True, radius_adj=0.14)
    text_box(s, lx + 0.42, by + 0.04, lw - 0.84, 0.62,
             "Knight Capital, 2012: $440M / ~45 мин на детерм. алгоритме "
             "— тот же класс ошибки",
             size=11.5, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    rx, rw = 7.80, 5.00
    ocean_box(s, rx, ly, rw, 1.62)
    text_box(s, rx + 0.26, ly + 0.14, rw - 0.52, 0.30,
             "Zillow vs Opendoor", size=14, bold=True, color=MID)
    text_box(s, rx + 0.26, ly + 0.48, rw - 0.52, 1.02,
             "Тот же тип ИИ, тот же период. Opendoor выжил за счёт более "
             "консервативного спреда и риск-дизайна.",
             size=12.5, color=DEEP, line_spacing=1.16)
    teal_callout(s, rx, ly + 1.78, rw, 2.32,
                 "Опасно, когда одновременно: (а) выход запускает крупное "
                 "необратимое действие авто, (б) среда нестационарна, "
                 "(в) нет circuit-breaker.\n\nАльтернатива: узкий сегмент + "
                 "human-gate + live error-monitoring + circuit-breaker.",
                 size=12, bold=False)
    gold_callout(s, 0.55, 5.40, 12.25, 0.90,
                 "Модель делала ровно то, для чего предназначена. Ошибка — "
                 "инженерное решение, КУДА подключить её выход. Тот же ИИ, "
                 "другое суждение → банкротство vs выживание.", size=14)
    speaker_notes(s, load_notes("s09"))


def build_s10(p):
    build_section_divider(
        p, 2, "Поиск аномалий:\nфрод и AML",
        "Прогноз — про будущее; теперь про настоящее — поймать аномалию за "
        "миллисекунды, пока платёж ещё не прошёл.", "s10")


def build_s11(p):
    """assertion_visual — anomaly cloud diagram + why-not."""
    s = blank(p)
    slide_title(s, "Учить «норму клиента» и ловить отклонение — а не учить "
                   "«как выглядит фрод».", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.20, 5.55, 4.05
    ocean_box(s, lx, ly, lw, lh)
    add_image(s, DIAGRAMS / "d11-anomaly-cloud.png",
              lx + 0.20, ly + 0.18, lw - 0.40, lh - 0.36)
    rx, rw = 6.30, 6.50
    ocean_box(s, rx, ly, rw, 4.05)
    text_box(s, rx + 0.26, ly + 0.16, rw - 0.52, 0.62,
             "Задача: поймать мошенническую транзакцию в реальном времени, "
             "где почти всё легитимно, а доля фрода крайне мала.",
             size=13, bold=True, color=MID, line_spacing=1.14)
    text_box(s, rx + 0.26, ly + 0.86, rw - 0.52, 0.56,
             "Тип ИИ — поиск аномалий (anomaly detection): модель строит "
             "«норму» клиента и сигналит об отклонениях.",
             size=12.5, color=DEEP, line_spacing=1.14)
    whys = [
        ("не классификация", "фрод редок и постоянно меняет форму "
                              "(движущаяся цель)"),
        ("не прогноз ряда", "вопрос — не «следующее значение», а «насколько "
                            "не похоже на норму»"),
        ("не LLM", "транзакция структурирована; задача геометрическая, не "
                   "языковая"),
    ]
    fy = ly + 1.52
    for a, b in whys:
        text_box(s, rx + 0.26, fy, 1.85, 0.56, a, size=12.5, bold=True,
                 color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, rx + 2.16, fy, rw - 2.42, 0.56, b, size=12, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        fy += 0.60
    filled_rect(s, rx + 0.26, fy + 0.06, rw - 0.52, 0.62, SURFACE,
                stroke=SOFT_GREY, stroke_pt=0.75, radius=True,
                radius_adj=0.12)
    text_box(s, rx + 0.40, fy + 0.06, rw - 0.80, 0.62,
             "AML — подмножество той же задачи; жёсткие законодательные "
             "пороги = детерминированные правила, не модель.",
             size=11, italic=True, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "Модель не «знает», что это фрод — она знает, что это НЕ "
                 "ПОХОЖЕ НА НОРМУ, и поднимает флаг для проверки.", size=14)
    speaker_notes(s, load_notes("s11"))


def build_s12(p):
    """case_study — fraud chart + hint."""
    s = blank(p)
    slide_title(s, "Антифрод — это «норма / отклонение» в реальном времени, "
                   "не генерация.", size=23, y=0.34, h=0.58, w=12.25)
    # 3 SEPARATE stat-plates — each its own unit (no false shared axis)
    text_box(s, 0.55, 1.16, 12.25, 0.30,
             "Что даёт поиск аномалий в антифроде — каждая метрика в своих "
             "единицах:", size=13.5, bold=True, color=MID)
    plates = [
        ("Stripe Radar", "−32% фрода", "при одобрении легитимных > 99%",
         False),
        ("JPMorgan", "−30% ложных", "меньше ложных срабатываний",
         False),
        ("Visa · предотвращено", "~$40 млрд", "мошеннических операций "
         "(FY2023)", True),
    ]
    n = 3
    gap = 0.24
    cw = (12.25 - gap * (n - 1)) / n
    px, py, ph = 0.55, 1.54, 2.16
    for lab, val, sub, hi in plates:
        bg = GOLD_TINT if hi else SURFACE
        edge = GOLD if hi else LIGHT
        ocean_box(s, px, py, cw, ph, fill=bg, stroke=edge,
                  stroke_pt=(2.0 if hi else 1.5))
        text_box(s, px + 0.20, py + 0.20, cw - 0.40, 0.34, lab,
                 size=14, bold=True, color=MID, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, px + 0.16, py + 0.62, cw - 0.32, 0.78, val,
                 size=33, bold=True, color=(GOLD if hi else DEEP),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, px + 0.20, py + 1.46, cw - 0.40, 0.58, sub,
                 size=12, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.10)
        px += cw + gap
    # Russia context + the metric hint, side by side below plates
    lx2, lw2 = 0.55, 6.00
    ocean_box(s, lx2, 3.92, lw2, 1.04)
    text_box(s, lx2 + 0.24, 4.04, lw2 - 0.48, 0.28, "Россия",
             size=13, bold=True, color=MID)
    text_box(s, lx2 + 0.24, 4.34, lw2 - 0.48, 0.56,
             "Поиск аномалий — стандартная практика крупных банков "
             "(по материалам Банка России, 20.11.2025).",
             size=11.5, color=DEEP, line_spacing=1.12)
    rx2, rw2 = 6.80, 6.00
    teal_callout(s, rx2, 3.92, rw2, 1.04,
                 "Заметьте: «снижение ЛОЖНЫХ срабатываний» — ключевая "
                 "метрика антифрода НЕ «точность вообще», а соотношение "
                 "двух типов ошибок. Разберём дальше.",
                 size=11.5, bold=False)
    gold_callout(s, 0.55, 5.12, 12.25, 0.86,
                 "Разные единицы — НЕ одна шкала: % снижения, % ложных, "
                 "млрд $. И всё равно главная метрика — не «accuracy», а "
                 "FP / FN раздельно.", size=14)
    footer(s, "Stripe/JPMorgan/Visa — заявленные компаниями; Visa — по "
              "сообщениям Reuters/CNBC, июль 2024; РФ — по материалам Банка "
              "России, 20.11.2025; числа сверяются в день лекции.")
    speaker_notes(s, load_notes("s12"))


def build_s13(p):
    """comparison / schema_matrix — confusion matrix 2x2 + accuracy lie."""
    s = blank(p)
    slide_title(s, "«Точность 99,9%» при сильном дисбалансе обманчива.",
                size=24, y=0.34, h=0.58, w=12.25)
    bx, by, bw, bh = 0.55, 1.18, 7.55, 3.10
    ocean_box(s, bx, by, bw, bh)
    # 2x2 matrix
    mx, my = bx + 0.30, by + 0.24
    col_lab = ["Было фродом", "Было легитимным"]
    cw2, ch2 = 3.30, 1.10
    text_box(s, mx + 1.10, my - 0.02, cw2, 0.26, col_lab[0],
             size=11.5, bold=True, color=MID, align=PP_ALIGN.CENTER)
    text_box(s, mx + 1.10 + cw2, my - 0.02, cw2, 0.26, col_lab[1],
             size=11.5, bold=True, color=MID, align=PP_ALIGN.CENTER)
    rlab = ["Система:\nфрод", "Система:\nпропуск"]
    cells = [
        [("TP", "поймали (хорошо)", TEAL_TINT, TEAL),
         ("FP", "честный клиент заблокирован\n(ошибка первого рода)",
          GOLD_TINT, GOLD)],
        [("FN", "деньги ушли мошеннику\n(ошибка второго рода)",
          RGBColor(0xE6, 0xE9, 0xF2), DEEP),
         ("TN", "корректно пропустили\n(хорошо)", TEAL_TINT, TEAL)],
    ]
    for ri in range(2):
        text_box(s, mx - 0.04, my + 0.30 + ri * (ch2 + 0.12), 1.05, ch2,
                 rlab[ri], size=11, bold=True, color=MID,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                 line_spacing=1.0)
        for ci in range(2):
            tag, desc, fill, fg = cells[ri][ci]
            xx = mx + 1.10 + ci * (cw2 + 0.05)
            yy = my + 0.30 + ri * (ch2 + 0.12)
            filled_rect(s, xx, yy, cw2, ch2, fill,
                        stroke=fg, stroke_pt=1.2, radius=True,
                        radius_adj=0.08)
            text_box(s, xx + 0.12, yy + 0.08, cw2 - 0.22, 0.32, tag,
                     size=16, bold=True, color=fg)
            text_box(s, xx + 0.12, yy + 0.44, cw2 - 0.22, ch2 - 0.52, desc,
                     size=11, color=DEEP, line_spacing=1.05)
    rx, rw = 8.30, 4.50
    ocean_box(s, rx, by, rw, 3.10)
    text_box(s, rx + 0.24, by + 0.16, rw - 0.48, 0.30,
             "Почему accuracy лжёт", size=13.5, bold=True, color=MID)
    text_box(s, rx + 0.24, by + 0.52, rw - 0.48, 1.10,
             "Поток 1 000 000, фрода 1000 (0,1%). Модель, говорящая «всё "
             "легитимно», даёт accuracy 99,9% — и совершенно бесполезна.",
             size=12, color=DEEP, line_spacing=1.18)
    filled_rect(s, rx + 0.24, by + 1.66, rw - 0.48, 0.02, SOFT_GREY)
    text_box(s, rx + 0.24, by + 1.78, rw - 0.48, 0.30,
             "cost-sensitive", size=13.5, bold=True, color=MID)
    text_box(s, rx + 0.24, by + 2.12, rw - 0.48, 0.90,
             "FP и FN различны по ЦЕНЕ (деньги, доверие). Минимизировать "
             "суммарную ожидаемую стоимость, не «красивую» accuracy.",
             size=12, color=DEEP, line_spacing=1.18)
    text_box(s, 0.55, 4.40, 12.25, 0.26,
             "Формальный аппарат (sensitivity / specificity) строится в "
             "Лекции 7 на медицинском примере; здесь — рабочая интуиция.",
             size=11, italic=True, color=SLATE)
    gold_callout(s, 0.55, 4.78, 12.25, 0.80,
                 "Первый вопрос к любой цифре антифрода — НЕ «какая "
                 "точность», а «какие FP и FN раздельно и по какой цене».",
                 size=14)
    speaker_notes(s, load_notes("s13"))


def build_s14(p):
    """case_study — FP cost contrast + precision/recall + Knight 2."""
    s = blank(p)
    slide_title(s, "«Снизили FP на 25%» усредняет копеечный FP и "
                   "катастрофический FP.", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 7.05, 4.10
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.30,
             "Тот же класс ошибки — цена на порядки разная",
             size=13.5, bold=True, color=MID)
    fp = [
        ("FP №1 · $5 за кофе", "клиент повторил, забыл через час → "
                               "цена ≈ 0, ОБРАТИМО", False),
        ("FP №2 · $5000 за лечение за границей",
         "дозвониться трудно, время критично → НЕОБРАТИМО по последствиям",
         True),
    ]
    fy = ly + 0.52
    for ttl, body, hi in fp:
        bg = GOLD_TINT if hi else SURFACE
        filled_rect(s, lx + 0.26, fy, lw - 0.52, 0.86, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.42, fy + 0.08, lw - 0.84, 0.30, ttl,
                 size=12.5, bold=True, color=(DEEP if hi else MID))
        text_box(s, lx + 0.42, fy + 0.40, lw - 0.84, 0.42, body,
                 size=11, color=DEEP, line_spacing=1.05)
        fy += 0.94
    text_box(s, lx + 0.26, fy + 0.04, lw - 0.52, 0.74,
             "precision ↔ recall — КОМПРОМИСС, а не задача «обе в ноль». "
             "Лучшая модель сдвигает кривую, точку выбирает инженер по "
             "цене ошибки.",
             size=11.5, bold=True, color=TEAL, line_spacing=1.12)
    filled_rect(s, lx + 0.26, fy + 0.84, lw - 0.52, 0.46, TEAL_TINT,
                stroke=TEAL, stroke_pt=1.2, radius=True, radius_adj=0.16)
    text_box(s, lx + 0.40, fy + 0.84, lw - 0.80, 0.46,
             "Авто-блок необратимого без гейта — класс Zillow / Knight",
             size=11, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    rx, rw = 7.80, 5.00
    teal_callout(s, rx, ly, rw, 4.10,
                 "Критерий:\n\n• авто-жёсткий блок — только обратимое / "
                 "мелкое\n\n• крупное / необратимое → мягкий challenge "
                 "(3DS, звонок, пуш) + быстрый человеческий канал "
                 "разблокировки\n\n• жёсткий AML-порог — rules-engine "
                 "(закон, не «с вероятностью 0,97»)\n\nЦена FP усредняет "
                 "копеечный и катастрофический — измерять раздельно.",
                 size=13, bold=False)
    gold_callout(s, 0.55, 5.40, 12.25, 0.90,
                 "Масштаб не уменьшает вред: 0,5% ложных от миллиардов = "
                 "десятки миллионов заблокированных легитимных операций. За "
                 "каждой — конкретный человек.", size=14)
    speaker_notes(s, load_notes("s14"))


def build_s15(p):
    build_section_divider(
        p, 3, "Кредитный скоринг:\nпочему НЕ нейросеть",
        "Аномалию нашли; теперь решение, которое меняет жизнь клиента, — и "
        "почему здесь нельзя чёрный ящик.", "s15")


def build_s16(p):
    """assertion_visual — 4 argument cards + инспектор analogy."""
    s = blank(p)
    slide_title(s, "Чёрный ящик в скоринге неприменим в принципе — не "
                   "«менее удобен».", size=23, y=0.34, h=0.58, w=12.25)
    cards = [
        ("gavel", "1. Объяснимость = закон",
         "клиенту с отказом по закону положено объяснение (reason codes: "
         "«высокая долговая нагрузка»). «Алгоритм так решил» — недопустимо.",
         "mid", False),
        ("database", "2. Данные табличные",
         "глубокие сети сильны на тексте/звуке/изображении; на табличных "
         "бустинг конкурентен при несравнимо большей интерпретируемости.",
         "mid", False),
        ("badge-check", "3. Аудит > +1% точности",
         "регулятор должен воспроизвести решение, проверить отсутствие "
         "дискриминации. Стабильная объяснимая важнее непрозрачной точнее.",
         "mid", False),
        ("triangle-alert", "4. Что сломается с чёрным ящиком",
         "нельзя объяснить → нарушение · нельзя аудировать на bias → "
         "кризис (Apple Card) · нестабильно → не защитить перед "
         "регулятором.", "gold", True),
    ]
    # left: 4 argument cards (compact, single column band) ; right: analogy
    cw2, ch2 = 5.95, 0.96
    for i, (ic, ttl, body, var, hi) in enumerate(cards):
        cx = 0.55
        cy = 1.16 + i * (ch2 + 0.14)
        if hi:
            ocean_box(s, cx, cy, cw2, ch2, fill=GOLD_TINT, stroke=GOLD,
                      stroke_pt=2.0)
        else:
            ocean_box(s, cx, cy, cw2, ch2)
        icon(s, ic, cx + 0.18, cy + (ch2 - 0.36) / 2, 0.36, var)
        text_box(s, cx + 0.66, cy + 0.10, cw2 - 0.84, 0.30, ttl,
                 size=12.5, bold=True, color=(DEEP if hi else MID))
        text_box(s, cx + 0.66, cy + 0.40, cw2 - 0.84, ch2 - 0.48, body,
                 size=10.5, color=DEEP, line_spacing=1.08)
    # right — inspector analogy diagram (chapter §3.3, derived)
    rx, ry, rw, rh = 6.70, 1.16, 6.10, 4.54
    ocean_box(s, rx, ry, rw, rh)
    text_box(s, rx + 0.24, ry + 0.14, rw - 0.48, 0.30,
             "Аналогия-якорь: кредитный инспектор", size=13, bold=True,
             color=MID)
    add_image(s, DIAGRAMS / "d16-inspector-reason-codes.png",
              rx + 0.18, ry + 0.50, rw - 0.36, rh - 1.30)
    text_box(s, rx + 0.24, ry + rh - 0.74, rw - 0.48, 0.62,
             "Интерпретируемая модель показывает расчёт построчно (reason "
             "codes); чёрный ящик отказывается объяснять.",
             size=10.5, italic=True, color=SLATE, line_spacing=1.10)
    gold_callout(s, 0.55, 5.86, 12.25, 0.92,
                 "«Новое = нейросеть, значит лучше» в скоринге — структурная "
                 "ошибка: объяснимость здесь не бонус, а УСЛОВИЕ "
                 "ЗАКОННОСТИ.", size=14)
    speaker_notes(s, load_notes("s16"))


def build_s17(p):
    """case_study — Сбер stats + clarification."""
    s = blank(p)
    slide_title(s, "«100% решений ИИ» в РФ ≠ нейросеть-чёрный-ящик "
                   "бесконтрольно.", size=23, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.30, 5.55, 3.95
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.32,
             "Сбербанк · с февраля–марта 2024", size=13.5, bold=True,
             color=MID)
    sb = [("Решения по физлицам", "~100% ИИ", False),
          ("Скоринговая модель", "до 5000 параметров", False),
          ("Эффект ИИ по всем направлениям, 2023", "+350 млрд ₽", True)]
    cyy = ly + 0.56
    for lab, val, hi in sb:
        bg = GOLD_TINT if hi else SURFACE
        filled_rect(s, lx + 0.26, cyy, lw - 0.52, 1.00, bg,
                    stroke=(GOLD if hi else SOFT_GREY),
                    stroke_pt=(1.5 if hi else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, lx + 0.42, cyy + 0.10, lw - 0.84, 0.42, lab,
                 size=12, bold=True, color=MID, line_spacing=1.05)
        text_box(s, lx + 0.42, cyy + 0.50, lw - 0.84, 0.44, val,
                 size=20, bold=True, color=(GOLD if hi else DEEP))
        cyy += 1.08
    rx, rw = 6.30, 6.50
    ocean_box(s, rx, ly, rw, 3.95)
    text_box(s, rx + 0.26, ly + 0.14, rw - 0.52, 0.28,
             "«100% ИИ» означает на самом деле:", size=14, bold=True,
             color=MID)
    items = [
        ("высокая автоматизация конвейера на интерпретируемых моделях",
         False),
        ("reason codes на каждое решение", False),
        ("регуляторный надзор Банка России", False),
        ("сохранённое право клиента на сотрудника: > 80% "
         "финорганизаций дают opt-out на человека", True),
    ]
    iy = ly + 0.48
    for it, hi in items:
        ih = 0.74 if hi else 0.52
        if hi:
            filled_rect(s, rx + 0.20, iy, rw - 0.40, ih, GOLD_TINT,
                        stroke=GOLD, stroke_pt=1.5, radius=True,
                        radius_adj=0.12)
            circle(s, rx + 0.38, iy + ih / 2 - 0.06, 0.13, GOLD)
            text_box(s, rx + 0.66, iy, rw - 0.92, ih, it,
                     size=12, bold=True, color=DEEP,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
        else:
            circle(s, rx + 0.30, iy + 0.13, 0.12, MID)
            text_box(s, rx + 0.56, iy, rw - 0.82, ih, it,
                     size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=1.10)
        iy += ih + 0.14
    teal_callout(s, rx, ly + 3.24, rw, 0.60,
                 "Тип ИИ — табличный ML, выбранный из-за объяснимости, а не "
                 "вопреки. Тезис лекции на масштабе крупнейшего банка.",
                 size=11.5, bold=False)
    footer(s, "По заявлениям банка (TAdviser/AdIndex/Интерфакс, 2024) и "
              "материалам Банка России (20.11.2025); числа сверяются в "
              "день лекции.")
    speaker_notes(s, load_notes("s17"))


def build_s18(p):
    """case_study — Apple Card precise outcome + proxy-bias from scratch."""
    s = blank(p)
    slide_title(s, "Даже формально-законная модель без объяснимости создаёт "
                   "кризис.", size=23, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 6.35, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.32,
             "Apple Card / Goldman Sachs · 2019–2021", size=13.5,
             bold=True, color=MID)
    timeline = [
        "нояб. 2019: вирусный тред — лимит ×20 выше у мужа (рейтинг "
        "супруги выше)",
        "NYDFS открыл расследование (~400 000 заявителей)",
        "март 2021: NYDFS НЕ нашёл нарушения закона",
        "НО: клиенты не получали объяснений → регуляторно-репутационный "
        "кризис",
    ]
    ty = ly + 0.54
    for i, t in enumerate(timeline):
        hi = (i == 2)
        circle(s, lx + 0.26, ty + 0.06, 0.13, GOLD if hi else MID)
        text_box(s, lx + 0.52, ty, lw - 0.78, 0.74, t,
                 size=11.5, bold=hi, color=(DEEP if not hi else DEEP),
                 line_spacing=1.10)
        ty += 0.80
    text_box(s, lx + 0.26, ty + 0.02, lw - 0.52, 0.40,
             "Утверждать «доказана дискриминация» — фактологически "
             "неверно.", size=11, italic=True, color=SLATE,
             line_spacing=1.05)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 2.32)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.30,
             "Прокси-предвзятость «на пальцах»", size=13.5, bold=True,
             color=MID)
    text_box(s, rx + 0.24, ly + 0.48, rw - 0.48, 1.72,
             "Инженер не подаёт пол/расу. Но остаются признаки, "
             "КОРРЕЛИРУЮЩИЕ с ними: почтовый индекс, история трат, "
             "занятость. Модель, оптимизируя точность, КОСВЕННО "
             "восстанавливает запрещённый признак через прокси — без "
             "всякого умысла.",
             size=12, color=DEEP, line_spacing=1.18)
    teal_callout(s, rx, ly + 2.48, rw, 1.10,
                 "Критерий: регулируемое решение требует одновременно "
                 "reason codes + appeal path (человек на оспоренном) + "
                 "bias-аудит исходов ДО прода.", size=12, bold=False)
    text_box(s, rx + 0.04, ly + 3.66, rw, 0.36,
             "Канонический разбор механизма (Obermeyer/Optum) — Лекция 7.",
             size=10.5, italic=True, color=SLATE)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "«Мы не используем защищённые признаки» — НЕОБХОДИМОЕ, но "
                 "совершенно недостаточное условие. Единственное "
                 "доказательство — прямой аудит исходов по группам.",
                 size=14)
    speaker_notes(s, load_notes("s18"))


def build_s19(p):
    """case_study — criterion synthesis (automation without gate)."""
    s = blank(p)
    slide_title(s, "Чем выше автономия, тем строже должен быть гейт вокруг "
                   "неё.", size=23, y=0.34, h=0.58, w=12.25)
    bx, by, bw, bh = 0.55, 1.20, 12.25, 1.60
    ocean_box(s, bx, by, bw, bh)
    text_box(s, bx + 0.26, by + 0.14, bw - 0.52, 0.30,
             "Один класс ошибки — три разных типа ИИ:", size=14,
             bold=True, color=MID)
    three = [
        ("Zillow", "прогноз → авто необратимая скупка"),
        ("fraud-автоблок", "аномалия → авто необратимая блокировка"),
        ("Knight Capital", "детерм. алгоритм → авто необратимые заявки"),
    ]
    n = 3
    gap = 0.20
    cwt = (bw - 0.52 - gap * (n - 1)) / n
    tx = bx + 0.26
    for ttl, body in three:
        filled_rect(s, tx, by + 0.52, cwt, 0.94, SURFACE, stroke=SOFT_GREY,
                    stroke_pt=1.0, radius=True, radius_adj=0.10)
        text_box(s, tx + 0.12, by + 0.60, cwt - 0.24, 0.30, ttl,
                 size=13, bold=True, color=TEAL)
        text_box(s, tx + 0.12, by + 0.90, cwt - 0.24, 0.50, body,
                 size=11, color=DEEP, line_spacing=1.05)
        tx += cwt + gap
    gold_callout(s, 0.55, 3.00, 12.25, 1.42,
                 "Автоматика, исполняющая необратимые финансовые действия в "
                 "открытом контуре без kill-switch (ручного «всё стоп»), без "
                 "лимитов на объём/позицию, без circuit-breaker (авто-стоп "
                 "при аномалии) и без верифицированного развёртывания, "
                 "превращает обычную ошибку модели в СКОРОСТЬ РАЗОРЕНИЯ.",
                 size=14)
    teal_callout(s, 0.55, 4.46, 12.25, 1.30,
                 "Скоринг в РФ ~100% автоматизирован → критерий применим "
                 "остро. «100% ИИ» приемлемо ТОЛЬКО в обвязке: reason codes "
                 "+ человеческий канал + bias-аудит + мониторинг дрейфа + "
                 "надзор. Уберите обвязку — это Apple Card на масштабе "
                 "банка.", size=13, bold=False)
    text_box(s, 0.55, 5.92, 12.25, 0.84,
             "Высокая автономия — не цель и не зло сама по себе. Цель — "
             "автономия, вокруг которой стоят оплаченные гейты, "
             "пропорциональные цене ошибки и необратимости действия.",
             size=13, italic=True, color=MID, line_spacing=1.15)
    speaker_notes(s, load_notes("s19"))


def build_s20(p):
    build_section_divider(
        p, 4, "LLM в финансах:\nассистенты и fact-checking",
        "До сих пор — намеренно НЕ LLM. Теперь LLM — и сразу про его "
        "границы в регулируемой отрасли.", "s20")


def build_s21(p):
    """assertion_visual — where LLM fits + fact-integrity base substitution."""
    s = blank(p)
    slide_title(s, "LLM хорош ровно там, где задача текстово-диалоговая — и "
                   "плох там, где нет.", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.20, 6.35, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.16, lw - 0.52, 0.30,
             "Где LLM — правильный тип в банке:", size=13.5, bold=True,
             color=MID)
    fits = [
        "клиентская поддержка — ответы на типовые вопросы",
        "голосовой ассистент — распознать запрос, провести по сценарию",
        "помощь сотруднику — суммаризировать обращение, подсказать "
        "регламент",
        "объяснение продукта простым языком",
    ]
    fy = ly + 0.52
    for t in fits:
        circle(s, lx + 0.26, fy + 0.05, 0.12, TEAL)
        text_box(s, lx + 0.52, fy, lw - 0.78, 0.56, t,
                 size=12, color=DEEP, line_spacing=1.10)
        fy += 0.58
    text_box(s, lx + 0.26, fy + 0.04, lw - 0.52, 0.34,
             "Это ровно то, подо что LLM спроектирована (в отличие от "
             "Разделов 1–3).", size=11, italic=True, color=SLATE,
             line_spacing=1.05)
    filled_rect(s, lx + 0.26, fy + 0.46, lw - 0.52, 0.78, SURFACE,
                stroke=SOFT_GREY, stroke_pt=0.75, radius=True,
                radius_adj=0.10)
    text_box(s, lx + 0.42, fy + 0.46, lw - 0.84, 0.78,
             "Масштаб (verified): чат-бот Т-Банка > 40% обращений; ~70% "
             "банков планировали голос к 2025.",
             size=11.5, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.08)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 4.05, fill=TEAL_TINT, stroke=TEAL,
              stroke_pt=2.0)
    text_box(s, rx + 0.26, ly + 0.18, rw - 0.52, 0.30,
             "Почему НЕ «> 90% обращений банков»", size=14, bold=True,
             color=TEAL)
    text_box(s, rx + 0.26, ly + 0.58, rw - 0.52, 1.74,
             "Цифра ~90% реальна — но это доля звонков голосового "
             "ассистента В КОЛЛ-ЦЕНТРЕ ОДНОГО БАНКА, не доля обращений "
             "банков в целом. Это ПОДМЕНА БАЗЫ (класс 5 fact-checking).",
             size=13.5, color=DEEP, line_spacing=1.22)
    text_box(s, rx + 0.26, ly + 2.46, rw - 0.52, 1.50,
             "Лекция учит fact-checking — строить assertion на цифре с "
             "подменённой базой значило бы нарушить ровно тот принцип, "
             "которому она учит. Точная скромная цифра ценнее громкой.",
             size=12.5, italic=True, color=MID, line_spacing=1.20)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "Симметрия с Разделами 1–3: «разная задача — разный тип "
                 "ИИ» означает не «LLM плох», а «LLM хорош ровно там, где "
                 "задача текстово-диалоговая, и плох там, где нет».",
                 size=14)
    speaker_notes(s, load_notes("s21"))


def build_s22(p):
    """assertion_visual — fact-checking criterion + 5 classes + grounding."""
    s = blank(p)
    slide_title(s, "LLM здесь не источник истины, а интерфейс к источнику "
                   "истины.", size=23, y=0.34, h=0.58, w=12.25)
    # TOP-LEFT — criterion (3 args, compact) ; TOP-RIGHT — 5 error classes
    lx, ly, lw, lh = 0.55, 1.12, 6.10, 2.42
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.24, ly + 0.12, lw - 0.48, 0.28,
             "Когда LLM НЕ источник истины — 3 аргумента", size=12.5,
             bold=True, color=MID)
    args = [
        ("механизм", "генерирует правдоподобное, не извлекает факт "
                      "(«правдоподобно» ≠ «верно»)"),
        ("цена", "неверная ставка/условие = нарушение, отвечает "
                 "организация"),
        ("что сломается", "свободная генерация о фактах = генератор "
                          "дезинформации с юр. ответственностью"),
    ]
    ay = ly + 0.46
    for ttl, body in args:
        text_box(s, lx + 0.24, ay, 1.55, 0.62, ttl,
                 size=11.5, bold=True, color=TEAL,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text_box(s, lx + 1.82, ay, lw - 2.06, 0.62, body,
                 size=10.5, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.08)
        ay += 0.64
    rx, rw = 6.80, 6.00
    ocean_box(s, rx, ly, rw, lh)
    text_box(s, rx + 0.24, ly + 0.12, rw - 0.48, 0.28,
             "Пять классов ошибок ИИ-утверждения", size=12.5, bold=True,
             color=MID)
    classes = [
        "1. галлюцинация факта → сверить с первоисточником",
        "2. устаревшие данные → проверить дату/актуальность",
        "3. прокси-bias в выводе → аудит исходов, не входов",
        "4. обман метрикой → FP/FN раздельно в деньгах",
        "5. подмена базы → запросить «доля чего от чего»",
    ]
    cy = ly + 0.44
    for c in classes:
        text_box(s, rx + 0.26, cy, rw - 0.50, 0.38, c,
                 size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
        cy += 0.38
    # BOTTOM BAND — grounding analogy as visual anchor (chapter §4.3)
    gx, gy, gw, gh = 0.55, 3.70, 12.25, 1.62
    ocean_box(s, gx, gy, gw, gh)
    text_box(s, gx + 0.26, gy + 0.12, gw - 0.52, 0.26,
             "Аналогия-якорь grounding: студент угадывает уверенным тоном "
             "vs сначала открывает справочник",
             size=12, bold=True, color=MID)
    add_image(s, DIAGRAMS / "d22-grounding-student.png",
              gx + 0.30, gy + 0.42, gw - 0.60, gh - 0.54)
    gold_callout(s, 0.55, 5.52, 12.25, 1.24,
                 "Альтернатива: фиксированный факт → детерминированная "
                 "выборка (grounded RAG), не генерация. На лекции — "
                 "распознать класс (Understand); проверить 5 заявлений по "
                 "первоисточникам — Семинар 5 (Apply).", size=14)
    speaker_notes(s, load_notes("s22"))


def build_s23(p):
    """case_study — Air Canada callback + Klarna arc."""
    s = blank(p)
    slide_title(s, "Правильный тип ИИ без grounding и человеческого выхода "
                   "= Air Canada / Klarna.", size=21, y=0.34, h=0.58,
                w=12.25)
    bx, by, bw, bh = 0.55, 1.18, 12.25, 3.30
    ocean_box(s, bx, by, bw, bh)
    col_w = 6.0
    panels = [
        ("Кейс E — Air Canada (callback)",
         ["чат-бот сообщил несуществующую политику возврата",
          "пассажир действовал на основании ответа",
          "фев. 2024: суд — компания ОТВЕЧАЕТ за инфо своего чат-бота, "
          "обязал выплатить компенсацию",
          "класс: галлюцинация финфакта = юр. ответственность"]),
        ("Кейс F — Klarna (дуга 2023 → 2025)",
         ["LLM-ассистент: ~2/3 обращений, ~11 мин → < 2 мин",
          "заявленная экономия ~$40 млн/год (Klarna, 2024)",
          "2024: подавалось как «ИИ заменяет поддержку»",
          "сер. 2025: CSAT↓ → вернулись к найму людей"]),
    ]
    for j, (ttl, items) in enumerate(panels):
        px = bx + 0.22 + j * (col_w + 0.10)
        filled_rect(s, px, by + 0.18, col_w - 0.20, 0.46,
                    MID if j == 0 else TEAL)
        text_box(s, px, by + 0.18, col_w - 0.20, 0.46, ttl,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        iy = by + 0.78
        for it in items:
            circle(s, px + 0.10, iy + 0.06, 0.11,
                   MID if j == 0 else TEAL)
            text_box(s, px + 0.32, iy, col_w - 0.56, 0.56, it,
                     size=11.5, color=DEEP, line_spacing=1.08)
            iy += 0.60
    teal_callout(s, 0.55, 4.62, 12.25, 0.80,
                 "Критерий: фиксированная политика/тариф/право → "
                 "детерминированная выборка/grounded, не свободная "
                 "генерация; и при любом сценарии гарантирован путь к "
                 "человеку.", size=13, bold=False)
    gold_callout(s, 0.55, 5.56, 12.25, 1.20,
                 "Устойчива не замена, а AUGMENTATION: ИИ на массовой "
                 "рутине + гарантированная человеческая эскалация на "
                 "эмоциональном / спорном / нестандартном хвосте.",
                 size=15)
    speaker_notes(s, load_notes("s23"))


def build_s24(p):
    """assertion_visual — PIVOT checkpoint Р4→Р5 (NOT retro-summary;
    distinct gold angle from s30 payoff; two-level insight = chapter §4.5).
    """
    s = blank(p)
    slide_title(s, "Чекпойнт: 4 типа из 5 пройдены — поворот к «человек + "
                   "обвязка».", size=22, y=0.34, h=0.58, w=12.25)
    # progress checkpoint — 4 done, 1 to go (explicit, not retro-summary)
    bx, by, bw, bh = 0.55, 1.18, 12.25, 1.46
    ocean_box(s, bx, by, bw, bh)
    text_box(s, bx + 0.26, by + 0.12, bw - 0.52, 0.28,
             "Прогресс по 5 типам — у 4 пройденных провал одной формы:",
             size=13.5, bold=True, color=MID)
    prog = [("✓ Прогноз", "Zillow", TEAL),
            ("✓ Аномалии", "fraud-FP", TEAL),
            ("✓ Скоринг", "Apple Card", TEAL),
            ("✓ LLM", "Air Canada", TEAL),
            ("→ Recsys", "следующий", GOLD)]
    n = 5
    gap = 0.16
    cwf = (bw - 0.52 - gap * (n - 1)) / n
    fx = bx + 0.26
    for ttl, body, col in prog:
        nxt = (col == GOLD)
        filled_rect(s, fx, by + 0.48, cwf, 0.82,
                    GOLD_TINT if nxt else SURFACE,
                    stroke=(GOLD if nxt else SOFT_GREY),
                    stroke_pt=(1.5 if nxt else 1.0), radius=True,
                    radius_adj=0.10)
        text_box(s, fx + 0.06, by + 0.54, cwf - 0.12, 0.30, ttl,
                 size=12.5, bold=True, color=col, align=PP_ALIGN.CENTER)
        text_box(s, fx + 0.06, by + 0.86, cwf - 0.12, 0.36, body,
                 size=10, color=DEEP, align=PP_ALIGN.CENTER,
                 line_spacing=1.0)
        fx += cwf + gap
    # two-level insight (chapter §4.5 — kept; it is the pivot's substance)
    teal_callout(s, 0.55, 2.80, 12.25, 1.16,
                 "Уровень 1 «какой тип ИИ» — решается структурой задачи, "
                 "чаще однозначно. Уровень 2 «что вокруг ИИ на цене ошибки» "
                 "— проектируется отдельно. Поворот лекции: дальше внимание "
                 "уходит с типа на обвязку.", size=13.5, bold=True)
    # the pivot itself — why the harmless-looking type is next
    filled_rect(s, 0.55, 4.12, 12.25, 1.16, SURFACE, stroke=LIGHT,
                stroke_pt=1.5, radius=True, radius_adj=0.07)
    icon(s, "shopping-cart", 0.85, 4.46, 0.48, "mid")
    text_box(s, 1.55, 4.22, 11.00, 0.98,
             "Последний тип — recsys / pricing — кажется самым "
             "БЕЗОБИДНЫМ: цена ошибки около нуля. Именно поэтому он "
             "опаснее остальных: провал ТИХИЙ — система рапортует успех "
             "по своей метрике.", size=13, color=DEEP, line_spacing=1.16,
             anchor=MSO_ANCHOR.MIDDLE)
    gold_callout(s, 0.55, 5.46, 12.25, 1.30,
                 "Вопрос на оставшийся раздел — НЕ «какой тип», а «почему "
                 "самый безобидный по цене ошибки тип может оказаться "
                 "опаснее громких провалов?»", size=15)
    speaker_notes(s, load_notes("s24"))


def build_s25(p):
    build_section_divider(
        p, 5, "Рекомендации и\nдинамическое ценообразование",
        "Последний тип — самый «безобидный» по цене ошибки, но с самыми "
        "тонкими патологиями.", "s25")


def build_s26(p):
    """comparison — collaborative vs content-based + user×item matrix."""
    s = blank(p)
    slide_title(s, "Два базовых подхода рекомендаций — и слабость у каждого "
                   "по имени.", size=23, y=0.34, h=0.58, w=12.25)
    # TOP — comparison table (full width, compact, parallel structure)
    bx, by, bw, bh = 0.55, 1.12, 12.25, 1.90
    ocean_box(s, bx, by, bw, bh)
    col_w = (bw - 0.40) / 2
    heads = ["Collaborative — «спроси похожих на тебя»",
             "Content-based — «ещё похожее по описанию»"]
    for j, hd in enumerate(heads):
        px = bx + 0.20 + j * col_w
        filled_rect(s, px, by + 0.16, col_w - 0.10, 0.42,
                    MID if j == 0 else TEAL)
        text_box(s, px, by + 0.16, col_w - 0.10, 0.42, hd,
                 size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("матрица пользователь × товар", "атрибуты товара (жанр, бренд)"),
        ("⊕ ловит неожиданные связи", "⊕ нет cold-start нового товара"),
        ("⊖ cold-start + popularity bias", "⊖ over-specialization (ниша)"),
    ]
    yy = by + 0.62
    for a, b in rows:
        for j, cc in enumerate((a, b)):
            px = bx + 0.20 + j * col_w
            filled_rect(s, px, yy, col_w - 0.10, 0.40, SURFACE,
                        stroke=SOFT_GREY, stroke_pt=0.75, radius=True,
                        radius_adj=0.14)
            text_box(s, px + 0.18, yy, col_w - 0.36, 0.40, cc,
                     size=11, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.LEFT, line_spacing=1.0)
        yy += 0.42
    # BOTTOM BAND — three sellers analogy as the main teaching visual
    gx, gy, gw, gh = 0.55, 3.16, 12.25, 1.78
    ocean_box(s, gx, gy, gw, gh)
    text_box(s, gx + 0.26, gy + 0.10, gw - 0.52, 0.26,
             "Аналогия-якорь: три продавца (collaborative · content · "
             "hybrid)", size=12, bold=True, color=MID)
    add_image(s, DIAGRAMS / "d26b-three-sellers.png",
              gx + 0.30, gy + 0.40, gw - 0.60, gh - 0.52)
    gold_callout(s, 0.55, 5.12, 12.25, 0.92,
                 "Запомнить: ЧТО это за подходы и КАКАЯ слабость по "
                 "названию у каждого — collaborative (паттерн поведения) vs "
                 "content (суть товара).", size=14)
    speaker_notes(s, load_notes("s26"))


def build_s27(p):
    """assertion_visual — hybrid + chart + filter bubble + dynamic pricing."""
    s = blank(p)
    slide_title(s, "Гибрид смягчает слабости — но без явного разнообразия "
                   "пузырь не отменяет.", size=22, y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 6.35, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.14, lw - 0.52, 0.28,
             "Hybrid recommender", size=14, bold=True, color=MID)
    text_box(s, lx + 0.26, ly + 0.46, lw - 0.52, 0.78,
             "collaborative + content-based + контекст. Контентная часть "
             "закрывает cold-start, коллаборативная ломает "
             "over-specialization. Большинство промышленных систем.",
             size=11.5, color=DEEP, line_spacing=1.14)
    # 2 stat-plates instead of ghost bars — readable + honest (estimate)
    pgap = 0.20
    pcw = (lw - 0.52 - pgap) / 2
    for i, (co, val, lab) in enumerate([
            ("Amazon", "~35%", "выручки — от рекомендаций"),
            ("Netflix", "~75%", "просмотров — от рекомендаций")]):
        ppx = lx + 0.26 + i * (pcw + pgap)
        filled_rect(s, ppx, ly + 1.34, pcw, 1.46, SURFACE, stroke=LIGHT,
                    stroke_pt=1.5, radius=True, radius_adj=0.10)
        text_box(s, ppx + 0.10, ly + 1.44, pcw - 0.20, 0.28, co,
                 size=12.5, bold=True, color=MID, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, ppx + 0.06, ly + 1.74, pcw - 0.12, 0.62, val,
                 size=30, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, ppx + 0.10, ly + 2.36, pcw - 0.20, 0.38, lab,
                 size=10.5, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, lx + 0.26, ly + 2.92, lw - 0.52, 0.98,
             "Историческая оценка (восходит к McKinsey ~2013), НЕ свежий "
             "верифицированный показатель — поэтому числом, не графиком-"
             "«фактом».",
             size=10.5, italic=True, color=SLATE, line_spacing=1.12)
    rx, rw = 7.10, 5.70
    ocean_box(s, rx, ly, rw, 1.92)
    text_box(s, rx + 0.24, ly + 0.14, rw - 0.48, 0.28,
             "Filter bubble (информационный пузырь)", size=13, bold=True,
             color=MID)
    text_box(s, rx + 0.24, ly + 0.46, rw - 0.48, 1.36,
             "показал похожее → пользователь взаимодействует с похожим → "
             "модель показывает ещё более похожее. Сужение разнообразия — "
             "НЕ злой умысел, а следствие оптимизации на краткосрочную "
             "релевантность.",
             size=11.5, color=DEEP, line_spacing=1.16)
    ocean_box(s, rx, ly + 2.08, rw, 1.97)
    text_box(s, rx + 0.24, ly + 2.22, rw - 0.48, 0.28,
             "Dynamic pricing", size=13, bold=True, color=MID)
    text_box(s, rx + 0.24, ly + 2.54, rw - 0.48, 1.42,
             "автоподстройка цены под спрос/время/контекст. Восприятие "
             "справедливости и закон — ОГРАНИЧЕНИЯ задачи, а НЕ "
             "оптимизируемые переменные.",
             size=11.5, color=DEEP, line_spacing=1.16)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "Filter bubble возникает без всякого умысла — прямое "
                 "следствие «угадать, что купишь сейчас». Гибрид это "
                 "сужение автоматически НЕ отменяет.", size=14)
    footer(s, "Amazon ~35% / Netflix ~75% — историческая оценка (McKinsey "
              "~2013), НЕ свежий показатель; Ozon/WB — точную долю компании "
              "не раскрывают.")
    speaker_notes(s, load_notes("s27"))


def build_s28(p):
    """case_study — proxy≠goal + Wendy's + criterion."""
    s = blank(p)
    slide_title(s, "Система рапортует успех по своей метрике — пока "
                   "разрушает, что в метрику не попало.", size=21,
                y=0.34, h=0.58, w=12.25)
    lx, ly, lw, lh = 0.55, 1.18, 7.05, 4.05
    ocean_box(s, lx, ly, lw, lh)
    text_box(s, lx + 0.26, ly + 0.14, lw - 0.52, 0.28,
             "Корень — прокси ≠ цель", size=14, bold=True, color=MID)
    text_box(s, lx + 0.26, ly + 0.46, lw - 0.52, 0.94,
             "система оптимизирует измеримое в моменте (ПРОКСИ: клики, "
             "время, маржа) ≠ настоящая цель (доверие, благополучие). "
             "Расхождение → filter bubble, dark patterns, гомогенизация.",
             size=11.5, color=DEEP, line_spacing=1.16)
    filled_rect(s, lx + 0.26, ly + 1.50, lw - 0.52, 0.34, GOLD_TINT,
                stroke=GOLD, stroke_pt=1.2, radius=True, radius_adj=0.2)
    text_box(s, lx + 0.40, ly + 1.50, lw - 0.80, 0.34,
             "Кейс G — Wendy's, февраль 2024", size=12.5, bold=True,
             color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, lx + 0.26, ly + 1.94, lw - 0.52, 1.96,
             "~$20 млн в цифровые меню-борды с динамическим "
             "ценообразованием → СМИ прочитали как surge-pricing → "
             "#BoycottWendys, конкурент обыграл в рекламе → откат за дни. "
             "Деньги вложены, технология доступна — провалилось "
             "игнорирование того, что справедливость восприятия = "
             "ОГРАНИЧЕНИЕ, а не переменная.",
             size=11.5, color=DEEP, line_spacing=1.18)
    rx, rw = 7.80, 5.00
    # counter-weight: «прокси ≠ цель» mini-diagram (chapter §5.5, derived)
    ocean_box(s, rx, ly, rw, 1.94)
    text_box(s, rx + 0.22, ly + 0.12, rw - 0.44, 0.26,
             "Почему провал ТИХИЙ:", size=12.5, bold=True, color=MID)
    # rising proxy line vs flat/declining true goal
    gx0, gy0, gx1 = rx + 0.30, ly + 1.62, rx + rw - 0.30
    connector(s, gx0, gy0, gx0, ly + 0.52, color=SLATE, width=1.2)
    connector(s, gx0, gy0, gx1, gy0, color=SLATE, width=1.2)
    connector(s, gx0, ly + 1.42, gx1, ly + 0.62, color=GOLD, width=3.0)
    connector(s, gx0, ly + 1.12, gx1, ly + 1.30, color=TEAL, width=3.0,
              dash="dash")
    text_box(s, gx1 - 1.70, ly + 0.50, 1.70, 0.26, "прокси-метрика ↑",
             size=10, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    text_box(s, gx1 - 1.70, ly + 1.30, 1.70, 0.26, "истинная цель →",
             size=10, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    text_box(s, rx + 0.22, ly + 1.66, rw - 0.44, 0.24,
             "дашборд зелёный, доверие падает", size=10,
             italic=True, color=SLATE)
    teal_callout(s, rx, ly + 2.10, rw, 1.95,
                 "Критерий:\n• прокси ≠ цель → serendipity + объяснимость "
                 "+ аудит на дискриминацию\n• цена — решение ЧЕЛОВЕКА В "
                 "ПРАВОВОЙ РАМКЕ, не выход оптимизатора", size=12,
                 bold=False)
    gold_callout(s, 0.55, 5.36, 12.25, 0.92,
                 "Единый класс через всю лекцию: accuracy-обман, прокси-"
                 "bias скоринга, Klarna, filter bubble — ПРОКСИ ≠ ЦЕЛЬ. "
                 "Чем мощнее оптимизатор, тем дороже расхождение. Провал "
                 "тихий.", size=14)
    speaker_notes(s, load_notes("s28"))


def build_s29(p):
    """matrix / schema_matrix — task × AI-type, 6 rows."""
    s = blank(p)
    slide_title(s, "Назвать тип под структуру задачи — и иногда правильный "
                   "тип это не ИИ.", size=22, y=0.32, h=0.56, w=12.25)
    bx, by, bw, bh = 0.40, 1.00, 12.55, 4.80
    ocean_box(s, bx, by, bw, bh)
    tx, ty = bx + 0.16, by + 0.14
    col_w = [3.05, 4.95, 4.23]
    headers = ["Задача", "Тип ИИ · почему именно он", "Типичный провал"]
    hh = 0.46
    cx = tx
    for j, hd in enumerate(headers):
        filled_rect(s, cx, ty, col_w[j], hh, MID)
        text_box(s, cx + 0.10, ty, col_w[j] - 0.18, hh, hd,
                 size=12.5, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE,
                 align=(PP_ALIGN.LEFT if j else PP_ALIGN.LEFT))
        cx += col_w[j]
    rows = [
        ("trending-up", "Прогноз спроса / cash-flow / отток",
         "Прогноз рядов (ARIMA/бустинг) — числа с трендом, не текст",
         "distribution shift × необратимое авто (Zillow)", False),
        ("radar", "Фрод real-time / AML",
         "Поиск аномалий + rules-engine — норма/отклонение",
         "FP на масштабе; accuracy лжёт при дисбалансе", False),
        ("scale", "Кредитный скоринг",
         "Табличный ML (логрег/GBM+SHAP) — объяснимость = закон",
         "прокси-bias + непрозрачность (Apple Card)", False),
        ("message-circle", "Поддержка / голос / объяснение",
         "LLM (grounded) — текстово-диалоговая задача",
         "галлюцинация финфакта = юр.ответственность (Air Canada)",
         False),
        ("shopping-cart", "Рекомендации / ценообразование",
         "Recsys; pricing — оптимизатор в рамке",
         "прокси≠цель: filter bubble (Wendy's)", False),
        ("circle-slash", "Детерм. регуляторная задача (жёсткий AML-порог)",
         "Обычный код / rules-engine — НЕ ИИ; закон ≠ вероятность",
         "ИИ добавил бы недетерминизм + поверхность ошибки", True),
    ]
    rh = 0.62
    yy = ty + hh
    for ic, task, aitype, fail, hi in rows:
        cx = tx
        cells = [task, aitype, fail]
        for j, cc in enumerate(cells):
            bg = GOLD_TINT if hi else (
                WHITE if (rows.index(
                    (ic, task, aitype, fail, hi)) % 2 == 0) else SURFACE)
            filled_rect(s, cx, yy, col_w[j], rh, bg,
                        stroke=(GOLD if hi else SOFT_GREY),
                        stroke_pt=(1.3 if hi else 0.75))
            if j == 0:
                icon(s, ic, cx + 0.10, yy + (rh - 0.30) / 2, 0.30,
                     "gold" if hi else "mid")
                text_box(s, cx + 0.48, yy, col_w[j] - 0.56, rh, cc,
                         size=10.5, bold=True,
                         color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                         line_spacing=1.0)
            else:
                text_box(s, cx + 0.12, yy, col_w[j] - 0.22, rh, cc,
                         size=10.5, bold=(hi and j == 1),
                         color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                         line_spacing=1.02)
            cx += col_w[j]
        yy += rh
    gold_callout(s, 0.55, 5.92, 12.25, 0.86,
                 "Нижняя строка — намеренно. Жёсткое верифицируемое "
                 "регуляторное правило → детерминированный код, НЕ ИИ: "
                 "ИИ добавил бы недетерминизм там, где нужны точность и "
                 "аудит.", size=14)
    speaker_notes(s, load_notes("s29"))


def build_s30(p):
    """assertion_visual — FINAL PAYOFF: one principle, distinct from s29
    matrix (operational apparatus) and s24 pivot. Dominant principle
    statement on top; 5 cases as compact evidence strip below."""
    s = blank(p)
    slide_title(s, "Содержательный ответ на центральный вопрос лекции.",
                size=23, y=0.34, h=0.56, w=12.25)
    # DOMINANT principle band — the single payoff statement
    bx, by, bw, bh = 0.55, 1.06, 12.25, 1.84
    ocean_box(s, bx, by, bw, bh, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, bx + 0.50, by + 0.20, bw - 1.00, bh - 0.40,
             "Ни в одном из пяти провалов исправление — это «лучшая "
             "модель». Во всех — лучшее СУЖДЕНИЕ о том, что стоит вокруг "
             "ИИ на цене ошибки.",
             size=21, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, line_spacing=1.20)
    # 5 cases as a compact horizontal evidence strip (case → fix)
    text_box(s, 0.55, 3.06, 12.25, 0.28,
             "Пять кейсов лекции — одно и то же лекарство (НЕ модель):",
             size=12.5, bold=True, color=MID)
    strip = [
        ("Zillow", "human-gate"),
        ("Apple Card", "объяснимость + аудит"),
        ("Air Canada", "grounding"),
        ("Klarna", "augmentation"),
        ("Wendy's", "правовая рамка"),
    ]
    n = 5
    gap = 0.18
    cw = (12.25 - gap * (n - 1)) / n
    sx, sy, sh = 0.55, 3.42, 1.46
    for case, fix in strip:
        ocean_box(s, sx, sy, cw, sh)
        text_box(s, sx + 0.08, sy + 0.16, cw - 0.16, 0.34, case,
                 size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, sx + cw / 2 - 0.10, sy + 0.54, 0.20, 0.26, "↓",
                 size=15, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
        filled_rect(s, sx + 0.14, sy + 0.82, cw - 0.28, 0.50, TEAL_TINT,
                    stroke=TEAL, stroke_pt=1.0, radius=True,
                    radius_adj=0.16)
        text_box(s, sx + 0.18, sy + 0.82, cw - 0.36, 0.50, fix,
                 size=10.5, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        sx += cw + gap
    gold_callout(s, 0.55, 5.06, 12.25, 1.00,
                 "Это и есть ответ: не «ИИ хорош/плох», а «назвать тип, "
                 "обосновать выбор и спроектировать обвязку под цену "
                 "ошибки».", size=14)
    speaker_notes(s, load_notes("s30"))


def build_s31(p):
    """assertion_visual — security 2 panels (А law / Б CV) — AGGREGATED."""
    s = blank(p)
    slide_title(s, "Финданные, PII (персональные данные) и биометрию нельзя "
                   "в публичный LLM.", size=21, y=0.32, h=0.56, w=12.25)
    # TWO clearly separated panels with header plates + thick divider
    lx, ly, lw, lh = 0.55, 1.06, 6.00, 2.78
    ocean_box(s, lx, ly, lw, lh)
    filled_rect(s, lx, ly, lw, 0.50, MID, radius=True, radius_adj=0.06)
    icon(s, "lock", lx + 0.22, ly + 0.07, 0.36, "white")
    text_box(s, lx + 0.70, ly, lw - 0.90, 0.50,
             "(А) Данные и закон: ФЗ-152 / PII", size=13, bold=True,
             color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    a_pts = [
        "финданные + PII + биометрия = чувствительные; ФЗ-152 — "
        "локализация ПДн граждан РФ, биометрия = строгий режим",
        "публичное облако: данные покидают контур · нет контроля "
        "retention · аудируемость падает",
        "критерий — по чувствительности данных и режиму, НЕ по силе "
        "модели",
    ]
    ay = ly + 0.62
    for t in a_pts:
        circle(s, lx + 0.24, ay + 0.05, 0.11, MID)
        text_box(s, lx + 0.48, ay, lw - 0.72, 0.70, t,
                 size=11, color=DEEP, line_spacing=1.10)
        ay += 0.72
    # vertical divider between the two independent sub-topics
    filled_rect(s, 6.62, ly, 0.06, lh, LIGHT)
    rx, rw = 6.80, 6.00
    ocean_box(s, rx, ly, rw, lh)
    filled_rect(s, rx, ly, rw, 0.50, TEAL, radius=True, radius_adj=0.06)
    icon(s, "scan-face", rx + 0.22, ly + 0.07, 0.36, "white")
    text_box(s, rx + 0.70, ly, rw - 0.90, 0.50,
             "(Б) CV-пласт: KYC, биометрия, скрытый труд", size=13,
             bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    b_pts = [
        "KYC = идентификация клиента; liveness — перед камерой живой "
        "человек, не фото/маска/дипфейк (строгий режим)",
        "Just Walk Out (Amazon) свёрнут 2024: «автономная» касса "
        "опиралась на > 1000 ревьюеров в Индии — скрытый труд",
        "Bias компьютерного зрения углубляется в Лекции 7.",
    ]
    by2 = ly + 0.62
    for t in b_pts:
        circle(s, rx + 0.24, by2 + 0.05, 0.11, TEAL)
        text_box(s, rx + 0.48, by2, rw - 0.72, 0.70, t,
                 size=11, color=DEEP, line_spacing=1.10)
        by2 += 0.72
    # GOLD ANCHOR BAND — biometrics irreversible (password vs face) — large
    gx, gy, gw, gh = 0.55, 3.96, 12.25, 1.98
    filled_rect(s, gx, gy, gw, gh, GOLD_TINT, stroke=GOLD, stroke_pt=2.5,
                radius=True, radius_adj=0.05)
    text_box(s, gx + 0.34, gy + 0.16, 3.95, gh - 0.32,
             "Биометрия необратимо\nскомпрометирована\nпри утечке — та же\n"
             "логика «необратимое →\nстрогий гейт», что\nпрошла через всю "
             "лекцию.",
             size=14, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.16)
    add_image(s, DIAGRAMS / "d31-password-vs-face.png",
              gx + 4.40, gy + 0.14, gw - 4.70, gh - 0.28)
    gold_callout(s, 0.55, 6.08, 12.25, 0.78,
                 "«Полностью автономно» в маркетинге — гипотеза для "
                 "проверки, не факт (Amazon оспаривал масштаб).",
                 size=14)
    speaker_notes(s, load_notes("s31"))


def build_s32(p):
    """qa_minimal — checklist + Семинар 5 bridge + Q&A."""
    s = blank(p)
    set_slide_bg(s, WHITE)
    bx, by, bw, bh = 0.55, 0.45, 12.25, 2.70
    ocean_box(s, bx, by, bw, bh)
    text_box(s, bx + 0.26, by + 0.14, bw - 0.52, 0.30,
             "Прежде чем применять ИИ к финансовой / ритейл-задаче:",
             size=14, bold=True, color=MID)
    checks = [
        "1. Какой тип ИИ диктует структура задачи?",
        "2. Почему НЕ LLM здесь? (≥1 структурная причина)",
        "3. Можно ли решить вообще без ИИ? (правило → код)",
        "4. Обратимо ли действие на выходе? (необратимое → gate)",
        "5. Как проверить факт? (класс ошибки + принцип)",
        "6. Регулируемо ли решение? (reason codes / appeal)",
        "7. Кто отвечает и где человек обязателен?",
        "8. Есть ли PII / финданные / биометрия? → не в LLM; ФЗ-152",
    ]
    n = len(checks)
    col_w = (bw - 0.52) / 2
    for i, c in enumerate(checks):
        col = i % 2
        row = i // 2
        cx = bx + 0.26 + col * col_w
        cy = by + 0.52 + row * 0.52
        text_box(s, cx, cy, col_w - 0.16, 0.48, c,
                 size=12, color=DEEP, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    ocean_box(s, bx, 3.30, bw, 1.30, fill=GOLD_TINT, stroke=GOLD,
              stroke_pt=2.0)
    text_box(s, bx + 0.30, 3.42, bw - 0.6, 0.34,
             "Задание — Семинар 5 «Проверка фактов AI на финансовых "
             "данных»", size=14, bold=True, color=DEEP)
    text_box(s, bx + 0.30, 3.78, bw - 0.6, 0.72,
             "Команды за ~30 мин сами проверяют 5 реальных заявлений по "
             "первоисточникам (Банк России / ВЦИОМ), маркируют "
             "верно/частично/ложь. Лекция учит распознавать класс "
             "(Understand); семинар — самостоятельно верифицировать "
             "(Apply).",
             size=12, color=DEEP, line_spacing=1.16)
    text_box(s, 0.55, 4.78, 12.25, 1.45, "Вопросы",
             size=92, bold=True, color=DEEP, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text_box(s, 0.55, 6.30, 12.25, 0.50,
             "Палитра типов — линза для всех отраслевых лекций дальше.  "
             "Спасибо за внимание.", size=15, italic=True, color=LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    speaker_notes(s, load_notes("s32"))


# ============================================================
# Main
# ============================================================
def main():
    spec = load_deck()
    if spec is not None:
        n = len(spec["slides"])
        print(f"deck spec OK — {n} slides (deck.yaml + deck-part2.yaml), "
              f"totals {spec['totals'].get('slides')}")
    p = setup_pres()
    builders = [
        build_s01, build_s02, build_s03, build_s04, build_s04a,
        build_s05, build_s06, build_s07, build_s08, build_s09,
        build_s10, build_s11, build_s12, build_s13, build_s14,
        build_s15, build_s16, build_s17, build_s18, build_s19,
        build_s20, build_s21, build_s22, build_s23, build_s24,
        build_s25, build_s26, build_s27, build_s28, build_s29,
        build_s30, build_s31, build_s32]
    assert len(builders) == 33, f"expected 33 builders, got {len(builders)}"
    for b in builders:
        b(p)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(OUT))
    print(f"saved {OUT} — {len(p.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
