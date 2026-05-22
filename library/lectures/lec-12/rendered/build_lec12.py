"""Full 39-slide build of Лекции 12 «AI в автоматизации производства и цифровые двойники».

Source-of-truth: deck.yaml v1 + chapter v3 multi-part (~30k слов) + slides/*.md.
Issue #133.

Palette LOCKED v3: Ocean Gradient (#21295C / #065A82 / #1C7293) + Teal (#028090)
secondary + Gold (#F0AB00) ≥1×/slide.
Visual motif: «Ocean rounded box» (radius 12, surface #F4F7FA, stroke #1C7293 1.5pt).
Canvas: 13.333" × 7.5" (16:9). Pacing per deck.yaml = 75 мин.

Lec-N-1 pattern compliance: match lec-11 (cover + lecture-map + 8 section dividers +
dedicated Q&A; top progress bar только on dividers + cover).

Build via: python3 build_lec12.py — generates lec-12.pptx.
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# === Palette (LOCKED v3) ===
DEEP    = RGBColor(0x21, 0x29, 0x5C)
MID     = RGBColor(0x06, 0x5A, 0x82)
LIGHT   = RGBColor(0x1C, 0x72, 0x93)
TEAL    = RGBColor(0x02, 0x80, 0x90)
SURFACE = RGBColor(0xF4, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xF0, 0xAB, 0x00)
SLATE   = RGBColor(0x6B, 0x76, 0x85)
GOLD_TINT = RGBColor(0xFE, 0xF5, 0xE0)
TEAL_TINT = RGBColor(0xE6, 0xF2, 0xF4)
SOFT_GREY = RGBColor(0xE5, 0xEA, 0xF0)
DARK_GREY = RGBColor(0x4A, 0x55, 0x6B)
ROADMAP = RGBColor(0xD9, 0xE2, 0xEC)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "lec-12.pptx"
FONT = "Arial"

# 8 sections roadmap (matching deck.yaml structure)
SECTIONS = ["Двойник 2026", "A0 Наблюдать", "A1 Советовать", "A2 Замыкать",
            "A3 Автоном.", "Где НЕ AI", "OT/IT", "РФ + Закр."]


def disable_shadow(shp):
    sppr = shp._element.spPr
    for el in sppr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
        sppr.remove(el)
    etree.SubElement(sppr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")


def text_box(slide, x, y, w, h, text, *, size=16, bold=False, italic=False,
             color=DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    try: p.line_spacing = line_spacing
    except: pass
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline(slide, x, y, w, h, lines, *, size=14, bold=False, color=DEEP,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
              line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = line_spacing
        except: pass
        if isinstance(line, tuple):
            txt, opts = line
        else:
            txt, opts = line, {}
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
        if "align" in opts:
            p.alignment = opts["align"]
    return tb


def rounded_box(slide, x, y, w, h, *, fill=SURFACE, stroke=LIGHT, stroke_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = stroke
    shp.line.width = Pt(stroke_w)
    disable_shadow(shp)
    return shp


def rectangle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def circle(slide, x, y, w, h, *, fill=MID, stroke=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if stroke:
        shp.line.color.rgb = stroke
    else:
        shp.line.fill.background()
    disable_shadow(shp)
    return shp


def arrow_right(slide, x, y, w, h, *, fill=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    disable_shadow(shp)
    return shp


def add_image(slide, path, x, y, w, h):
    if not path or not Path(path).exists():
        return None
    try:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                         width=Inches(w), height=Inches(h))
    except Exception as e:
        print(f"Image fail {path}: {e}")
        return None


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])


def setup_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def roadmap_bar(slide, current_section):
    """8-section roadmap bar at top of section dividers + cover.
    current_section=0 means «pre-DT» (cover); 1..8 are actual sections."""
    bar_y = 0.32
    bar_h = 0.30
    total_w = 12.33
    seg_w = total_w / 8
    for i, name in enumerate(SECTIONS):
        x = 0.5 + i * seg_w
        is_active = (i == current_section - 1)
        rectangle(slide, x, bar_y, seg_w - 0.03, bar_h,
                  fill=GOLD if is_active else ROADMAP)
        text_box(slide, x, bar_y + 0.02, seg_w - 0.03, bar_h - 0.04,
                 f"{i+1}. {name}", size=9, bold=is_active,
                 color=DEEP if is_active else SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, text, *, y=7.05):
    text_box(slide, 0.5, y, 12.33, 0.35, text,
             size=11, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)


def attribution(slide, text, *, x=0.5, y=6.95, w=12.33, color=SLATE):
    text_box(slide, x, y, w, 0.3, text, size=10, italic=True,
             color=color, align=PP_ALIGN.LEFT)


# ====================================================================
# SECTION 0 — Cover + keystone + lecture-map
# ====================================================================

def s01_hero(p):
    """s01 — hero hook: Hannover Messe robotic hand. ≥40% area (6.7×6.0 = 40.2%)."""
    slide = blank(p)
    set_bg(slide, WHITE)
    # Hero image left — 6.7 × 6.0 = 40.2% canvas (≥40% threshold)
    img = ASSETS / "screenshots" / "s01-hannover-messe.jpg"
    if img.exists():
        add_image(slide, img, 0.4, 0.5, 6.7, 6.0)
    else:
        rectangle(slide, 0.4, 0.5, 6.7, 6.0, fill=SOFT_GREY)
        text_box(slide, 0.4, 3.0, 6.7, 0.6, "[Hannover Messe hero]",
                 size=16, color=SLATE, align=PP_ALIGN.CENTER)
    attribution(slide, "Hannover Messe 2016 · робот-манипулятор · Wikimedia · CC-BY-SA",
                x=0.4, y=6.55, w=6.7)
    # Right column — assertion + scale
    multiline(slide, 7.3, 0.6, 5.7, 5.5, [
        ("Лекция 12", {"size": 16, "bold": True, "color": LIGHT}),
        ("", {"size": 6}),
        ("AI поднимается", {"size": 32, "bold": True, "color": DEEP}),
        ("по шкале автономии", {"size": 32, "bold": True, "color": DEEP}),
        ("", {"size": 12}),
    ], line_spacing=1.1)
    # Ladder visual right-bottom
    rounded_box(slide, 7.3, 3.6, 5.7, 1.8, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    multiline(slide, 7.5, 3.7, 5.3, 1.6, [
        ("A0 → A1 → A2 → A3", {"size": 28, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("", {"size": 4}),
        ("наблюдать  →  советовать  →  замыкать петлю  →  автономно",
         {"size": 11, "italic": True, "color": DEEP, "align": PP_ALIGN.CENTER}),
    ], line_spacing=1.15)
    # Anchor
    rounded_box(slide, 7.3, 5.6, 5.7, 1.0, fill=SURFACE, stroke=LIGHT, stroke_w=1.5)
    text_box(slide, 7.5, 5.7, 5.3, 0.8,
             "Цифровой двойник — единственный известный мост между ступенями",
             size=13, bold=True, color=MID, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    footer(slide, "Курс «Применение AI в инженерии»")


def s02_cover(p):
    """s02 — cover with lecture title + roadmap + LO."""
    slide = blank(p)
    set_bg(slide, WHITE)
    roadmap_bar(slide, current_section=0)
    # Large decorative «12» — wider box prevents stacking
    text_box(slide, 0.0, 1.4, 4.5, 5.6, "12",
             size=240, bold=True, color=ROADMAP, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    # Title
    multiline(slide, 4.7, 1.3, 8.3, 2.5, [
        ("Лекция 12", {"size": 20, "bold": True, "color": LIGHT}),
        ("", {"size": 4}),
        ("AI в автоматизации", {"size": 34, "bold": True, "color": DEEP}),
        ("производства и цифровые", {"size": 34, "bold": True, "color": DEEP}),
        ("двойники", {"size": 34, "bold": True, "color": DEEP}),
    ], line_spacing=1.05)
    # Meta
    multiline(slide, 4.7, 4.4, 8.3, 1.0, [
        ("Модуль 2 · Промышленный AI и автоматизация", {"size": 15, "color": LIGHT}),
        ("Студенты-инженеры 3 курса", {"size": 15, "color": LIGHT}),
    ])
    # Central question rounded box
    rounded_box(slide, 4.7, 5.5, 8.3, 1.5, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    multiline(slide, 4.9, 5.6, 7.9, 1.3, [
        ("Центральный вопрос:", {"size": 13, "bold": True, "color": MID}),
        ("На какой ступени автономии можно поднять процесс — A0/A1/A2/A3 —",
         {"size": 13, "color": DEEP}),
        ("и где двойник обязателен, а где AI не нужен вовсе?",
         {"size": 13, "color": DEEP}),
    ], line_spacing=1.3)
    footer(slide, "Курс «Применение AI в инженерии» · 2026")


def s03_lecture_map(p):
    """s03 — 8 horizontal section cards."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Маршрут лекции — восемь разделов",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.1, 12.33, 0.4,
             "От определения двойника к ступеням автономии A0–A3, к границам и российскому контексту",
             size=13, italic=True, color=LIGHT)
    sections_full = [
        ("1", "Двойник\n2026", "Kritzinger + ГОСТ + рынок + 75% провал", MID),
        ("2", "A0", "Vision QC + прогност. обслуж. + границы", LIGHT),
        ("3", "A1", "MES + предсказ. тревог + PLC Copilot", TEAL),
        ("4", "A2", "Yokogawa + двойник как песочница", MID),
        ("5", "A3", "Toyota Digit + 3 блокера", GOLD),
        ("6", "НЕ AI", "10 критериев + фарма+FDA", LIGHT),
        ("7", "OT/IT", "7 слоёв + Lighthouse", TEAL),
        ("8", "РФ\n+ Закр.", "ГОСТ + КАМАЗ + 4 роли", MID),
    ]
    card_w = 1.46
    gap = 0.04
    x0 = 0.6
    y = 1.7
    for i, (num, title, desc, accent) in enumerate(sections_full):
        x = x0 + i * (card_w + gap)
        rounded_box(slide, x, y, card_w, 5.1)
        circle(slide, x + card_w/2 - 0.4, y + 0.2, 0.8, 0.8, fill=accent)
        text_box(slide, x, y + 0.2, card_w, 0.8, num,
                 size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(slide, x + 0.05, y + 1.2, card_w - 0.1, 1.2, title,
                 size=15, bold=True, color=DEEP, align=PP_ALIGN.CENTER, line_spacing=1.1)
        text_box(slide, x + 0.1, y + 2.7, card_w - 0.2, 2.3, desc,
                 size=10, color=DARK_GREY, align=PP_ALIGN.CENTER, line_spacing=1.3)
    footer(slide, "Применимый инструмент для кармана — десять критериев «AI не подходит» на 6-м разделе")


def s04_keystone(p):
    """s04 — KEYSTONE: ladder A0→A3 + disclaimer ISA-95 + asymmetry."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Шкала автономии AI в производстве: A0 → A3",
             size=26, bold=True, color=DEEP)
    # First-line disambiguation
    rounded_box(slide, 0.5, 1.15, 12.33, 0.55, fill=SURFACE, stroke=LIGHT, stroke_w=1.0)
    text_box(slide, 0.7, 1.2, 12.0, 0.5,
             "Не путать с архитектурными слоями ISA-95 L0–L2: там — слои архитектуры; здесь — степени автономии AI",
             size=12, italic=True, color=MID, anchor=MSO_ANCHOR.MIDDLE)
    # Four-step ladder
    steps = [
        ("A0", "Наблюдать",   "Vision QC · прогностическое обслуживание",
         "AI выдаёт сигнал · решение принимает человек или другая система", LIGHT),
        ("A1", "Советовать",  "MES советующий · предсказание тревог · PLC Copilot",
         "AI предлагает действие · оператор даёт явное согласие", MID),
        ("A2", "Замыкать петлю", "Yokogawa FKDPP в JSR · энергооптимизация",
         "AI меняет параметры в безопасной зоне без согласия на каждое", TEAL),
        ("A3", "Автономно",    "Toyota Digit · BMW Leipzig humanoid (единицы)",
         "AI принимает решения без человека, c защитными ограничениями", GOLD),
    ]
    col_w = 2.95
    col_h = 4.4
    gap = 0.1
    x0 = 0.5
    y = 1.95
    for i, (lvl, name, examples, role, accent) in enumerate(steps):
        x = x0 + i * (col_w + gap)
        rounded_box(slide, x, y, col_w, col_h, fill=SURFACE, stroke=accent, stroke_w=2.0)
        # Level pill
        rectangle(slide, x, y, col_w, 0.7, fill=accent)
        text_box(slide, x + 0.1, y + 0.05, col_w - 0.2, 0.65, lvl,
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Name
        text_box(slide, x + 0.1, y + 0.85, col_w - 0.2, 0.55, name,
                 size=20, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        # Examples
        text_box(slide, x + 0.2, y + 1.55, col_w - 0.4, 1.4, examples,
                 size=12, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.35, italic=True)
        # Role
        rounded_box(slide, x + 0.2, y + 3.0, col_w - 0.4, 1.25, fill=WHITE, stroke=accent, stroke_w=1.0)
        text_box(slide, x + 0.3, y + 3.1, col_w - 0.6, 1.05, role,
                 size=11, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    # Visual bridge-arrow between A1 (i=1) and A2 (i=2) — gold motif "twin как мост"
    # x of A1 right edge = 0.5 + 1*(col_w+gap) + col_w = 0.5 + 3.05 + 2.95 = 6.5
    # x of A2 left edge = 0.5 + 2*(col_w+gap) = 0.5 + 6.1 = 6.6
    # Bridge arrow span: 6.45 (slightly overlapping A1) → 6.65 (slightly into A2), positioned mid-vertical
    bridge_y = y + 2.2  # middle of card height
    # Gold arc-bridge: small rounded box with "twin" label spanning gap
    rounded_box(slide, 6.35, bridge_y - 0.15, 0.8, 0.7, fill=GOLD, stroke=GOLD, stroke_w=2.0)
    text_box(slide, 6.35, bridge_y - 0.15, 0.8, 0.7, "DT",
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Two small arrows pointing from A1→DT and DT→A2
    arrow_right(slide, 6.15, bridge_y + 0.1, 0.18, 0.2, fill=GOLD)
    arrow_right(slide, 7.18, bridge_y + 0.1, 0.18, 0.2, fill=GOLD)
    text_box(slide, 5.85, bridge_y + 0.65, 1.8, 0.3, "двойник =\nмост",
             size=10, bold=True, italic=True, color=GOLD,
             align=PP_ALIGN.CENTER, line_spacing=1.1)
    # Bridge text
    rounded_box(slide, 0.5, 6.45, 12.33, 0.6, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    text_box(slide, 0.7, 6.5, 12.0, 0.55,
             "Цифровой двойник — мост между A1 и A2 · A3 в 2026 — единицы кейсов, основная масса A0–A2",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "Это несущий слайд лекции. Производственная автоматизация в 2026 — шкала автономии. A0: AI наблюдает, выдаёт сигнал, человек решает; Vision QC, прогностическое обслуживание. A1: AI советует, оператор соглашается; MES советующий, предсказание тревог, PLC Copilot. A2: AI меняет параметры процесса в безопасной зоне без согласия на каждое; Yokogawa FKDPP в JSR 35 дней, энергооптимизация. A3: AI действует автономно; Toyota Digit на RAV4, BMW Leipzig humanoid — единичные кейсы 2026. Якорь стандартов: SAE J3016 и ISO/IEC 22989 используют термин «уровни автономии». Уточнение: это шкала автономии AI, не архитектурные слои ISA-95 L0–L2 из лекции 11. Уточнение по асимметрии A3: единицы кейсов, основная масса A0–A2. Цифровой двойник — мост между A1 и A2: без двойника подъём = слепая вера, путь к Tesla 2018 и анонимному порту 2024.")


# ====================================================================
# SECTION DIVIDER GENERIC
# ====================================================================

def section_divider(p, num, current_section, title, subtitle, meta):
    """Generic section divider with large background number + roadmap bar."""
    slide = blank(p)
    set_bg(slide, WHITE)
    roadmap_bar(slide, current_section=current_section)
    # Large background number
    text_box(slide, 0.5, 1.2, 6.0, 5.5, str(num),
             size=380, bold=True, color=ROADMAP,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title right
    multiline(slide, 5.5, 2.4, 7.5, 4.0, [
        ("РАЗДЕЛ", {"size": 18, "bold": True, "color": LIGHT}),
        (title, {"size": 36, "bold": True, "color": DEEP}),
        ("", {"size": 8}),
        (subtitle, {"size": 15, "color": MID}),
        ("", {"size": 6}),
        (meta, {"size": 12, "italic": True, "color": GOLD}),
    ], line_spacing=1.15)
    add_notes(slide, f"Раздел {num}. {title}. {subtitle}. {meta}.")


# ====================================================================
# SECTION 1 — DT 2026 (§1) — s05 divider + s06-s10
# ====================================================================

def s05_div(p):
    section_divider(p, 1, 1, "Что такое цифровой двойник в 2026",
                    "Kritzinger taxonomy · ГОСТ Р 57700.37 · рынок $36→$180B · провал 75%",
                    "Раздел 1 · от определения к рыночной картине")


def s06_kritzinger(p):
    """s06 — Kritzinger taxonomy + ГОСТ."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Цифровой двойник — двусторонняя петля, а не CAD-картинка",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Таксономия Kritzinger 2018 + ГОСТ Р 57700.37-2021 — два эталонных определения",
             size=13, italic=True, color=LIGHT)
    # Three Kritzinger types — column cards
    types = [
        ("Цифровая модель", "(Digital Model)",
         "Нет потока данных от физического объекта",
         "Нет управляющего действия назад",
         "Пример: CAD-чертёж",
         LIGHT),
        ("Цифровая тень", "(Digital Shadow)",
         "Поток физика → цифра",
         "Нет управляющего действия назад",
         "Пример: панель мониторинга",
         MID),
        ("Цифровой двойник", "(Digital Twin)",
         "Двусторонняя петля физика ↔ цифра",
         "Управляющее действие назад через симуляцию",
         "Пример: Siemens Digital Twin Composer 2026",
         GOLD),
    ]
    col_w = 4.0
    gap = 0.13
    x0 = 0.5
    y = 1.75
    for i, (name_ru, name_en, flow, action, example, accent) in enumerate(types):
        x = x0 + i * (col_w + gap)
        rounded_box(slide, x, y, col_w, 4.4, fill=SURFACE, stroke=accent, stroke_w=2.0)
        rectangle(slide, x, y, col_w, 0.85, fill=accent)
        text_box(slide, x + 0.1, y + 0.08, col_w - 0.2, 0.45, name_ru,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text_box(slide, x + 0.1, y + 0.5, col_w - 0.2, 0.3, name_en,
                 size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
        multiline(slide, x + 0.2, y + 1.0, col_w - 0.4, 3.3, [
            ("Поток данных:", {"size": 11, "bold": True, "color": LIGHT}),
            (flow, {"size": 12, "color": DEEP}),
            ("", {"size": 6}),
            ("Управляющее действие:", {"size": 11, "bold": True, "color": LIGHT}),
            (action, {"size": 12, "color": DEEP}),
            ("", {"size": 6}),
            ("Где встречается:", {"size": 11, "bold": True, "color": LIGHT}),
            (example, {"size": 12, "italic": True, "color": MID}),
        ], line_spacing=1.3)
    # ГОСТ box below
    rounded_box(slide, 0.5, 6.3, 12.33, 0.75, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 0.7, 6.4, 12.0, 0.6,
             "ГОСТ Р 57700.37-2021 «Цифровые двойники изделий. Общие положения» — формальная регуляторная база РФ; синхронизирует таксономию Kritzinger",
             size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    add_notes(slide, "Kritzinger в 2018 году в IFAC-PapersOnLine ввёл три уровня. Цифровая модель: статичная репрезентация без потока данных. Цифровая тень: данные текут из физики в цифру, но управляющего действия назад нет — это панель мониторинга. Цифровой двойник: двусторонняя петля с возможностью симулировать управляющее действие назад. Только третий тип — настоящий двойник. Когда вендор показывает дашборд с 3D-визуализацией, спросите: где двусторонняя петля? Где возможность симулировать действие назад? ГОСТ Р 57700.37-2021 — формальная регуляторная база РФ, синхронизирует терминологию с международной таксономией.")


def s07_four_layers(p):
    """s07 — 4-layer architecture diagram. Hero reduced to 30-35% area; expanded layer cards."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Цифровой двойник требует четырёх слоёв одновременно",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Без любого из слоёв — это НЕ двойник, а только цифровая модель или тень",
             size=13, italic=True, color=LIGHT)
    # Hero image on left (reduced to ~3.6 wide × 5.0 high = 18 sq.in. ≈ 18% canvas; layers expand to ~65%)
    img = ASSETS / "screenshots" / "s07-siemens-amberg.jpg"
    if img.exists():
        add_image(slide, img, 0.5, 1.8, 3.6, 5.0)
        attribution(slide, "Siemens HQ Munich · Wikimedia · CC-BY-SA",
                    x=0.5, y=6.85, w=3.6)
    # Right: 4 layers stacked, expanded
    layers = [
        ("4. ИИ-потребители",
         "Vision QC · оптимизатор · RL-агент · интерфейс оператора",
         GOLD),
        ("3. Слой модели",
         "Физика + ML · симулятор оборудования · прокрутка времени",
         TEAL),
        ("2. Слой датчиков",
         "OPC UA + MQTT · частота опроса ≥10× полосы управления",
         MID),
        ("1. Физический актив",
         "Реактор · конвейер · линия · насос",
         LIGHT),
    ]
    x = 4.5
    y_start = 1.85
    layer_h = 1.20
    gap = 0.1
    for i, (name, desc, color) in enumerate(layers):
        y = y_start + i * (layer_h + gap)
        rounded_box(slide, x, y, 8.3, layer_h, fill=SURFACE, stroke=color, stroke_w=2.0)
        rectangle(slide, x, y, 0.5, layer_h, fill=color)
        text_box(slide, x + 0.7, y + 0.15, 7.5, 0.5, name,
                 size=17, bold=True, color=DEEP)
        text_box(slide, x + 0.7, y + 0.7, 7.5, 0.5, desc,
                 size=13, italic=True, color=DARK_GREY)
    # Arrows showing bidirectional flow
    text_box(slide, 4.25, 2.4, 0.3, 4.5, "↑↓\n↑↓\n↑↓\n↑↓",
             size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    footer(slide, "Без двусторонней петли (стрелки между слоями) — это не двойник, а только тень")


def s08_market(p):
    """s08 — Market chart."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.35, 12.33, 1.0,
             "Рынки растут на 35% в год — и одновременно 75% проектов проваливаются",
             size=22, bold=True, color=DEEP, line_spacing=1.1)
    text_box(slide, 0.5, 1.4, 12.33, 0.4,
             "Цифровые двойники $36→$180B · AI в производстве $155B · промышленный AI на OPC UA — уже $17B",
             size=13, italic=True, color=LIGHT)
    chart = ASSETS / "charts" / "s08-market.png"
    if chart.exists():
        rounded_box(slide, 0.5, 1.7, 8.0, 5.0)
        add_image(slide, chart, 0.7, 1.85, 7.6, 4.7)
    # Right column — failure stats
    rounded_box(slide, 8.7, 1.7, 4.1, 5.0, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    multiline(slide, 8.9, 1.85, 3.7, 4.8, [
        ("ПАРАЛЛЕЛЬНО:", {"size": 12, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("", {"size": 6}),
        ("75%", {"size": 48, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("проектов двойников не дают окупаемости",
         {"size": 12, "color": DEEP, "align": PP_ALIGN.CENTER}),
        ("", {"size": 10}),
        ("11%", {"size": 30, "bold": True, "color": MID, "align": PP_ALIGN.CENTER}),
        ("в нефтегазе дают ожидаемый эффект",
         {"size": 11, "color": DEEP, "align": PP_ALIGN.CENTER}),
        ("", {"size": 8}),
        ("14%", {"size": 30, "bold": True, "color": TEAL, "align": PP_ALIGN.CENTER}),
        ("пользователей: технология соответствует ожиданиям",
         {"size": 11, "color": DEEP, "align": PP_ALIGN.CENTER}),
    ], line_spacing=1.2)
    attribution(slide, "PatSnap / StartUs / Standard Bots / TheElec / EY / Build in Digital [41] 2024–2025", y=6.85)
    add_notes(slide, "Запомните четыре цифры. Рынок цифровых двойников: 36,19 миллиарда долларов в 2025 году, 180,28 миллиарда к 2030 году — среднегодовой темп роста 37,87%. Параллельный рынок AI в производстве: 155 миллиардов к 2030. Промышленный AI поверх OPC UA + MQTT — уже 17,15 миллиарда в 2026. И параллельно: 75% проектов двойников не дают окупаемости из-за слабого слоя данных; 11% в нефтегазе дают ожидаемый эффект; 14% пользователей считают, что технология соответствует ожиданиям. Эти числа описывают одну и ту же картину: гонка за технологией опережает инфраструктурную готовность.")


def s09_port_failure(p):
    """s09 — Southeast Asian Port failure case-card."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Цифровой двойник морского порта (анонимный кейс): $12 млн, 18 месяцев, списан 2024",
             size=24, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Двойник морского порта — 3D-визуализация без слоя данных = музей, а не двойник",
             size=13, italic=True, color=LIGHT)
    # Hero image
    img = ASSETS / "screenshots" / "s09-container-port.jpg"
    if img.exists():
        add_image(slide, img, 0.5, 1.75, 7.0, 5.0)
        attribution(slide, "Контейнерный порт Сингапура · Tanjong Pagar Terminal · Wikimedia · CC-BY-SA",
                    x=0.5, y=6.8, w=7.0)
    # Right column — failure analysis
    rounded_box(slide, 7.7, 1.75, 5.1, 5.0, fill=SURFACE, stroke=GOLD, stroke_w=2.0)
    rectangle(slide, 7.7, 1.75, 5.1, 0.7, fill=GOLD)
    text_box(slide, 7.7, 1.8, 5.1, 0.65, "ПРОВАЛ",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    multiline(slide, 7.9, 2.6, 4.7, 4.0, [
        ("Цифры:", {"size": 12, "bold": True, "color": MID}),
        ("· $12 миллионов — бюджет", {"size": 12, "color": DEEP}),
        ("· 18 месяцев — пилот", {"size": 12, "color": DEEP}),
        ("· 2024 — списан", {"size": 12, "color": DEEP}),
        ("", {"size": 8}),
        ("Три корневые причины:", {"size": 12, "bold": True, "color": MID}),
        ("1. Фрагментированные данные из 5 источников без единой схемы",
         {"size": 11, "color": DEEP}),
        ("2. Чрезмерный акцент на 3D-визуализации вместо потока данных",
         {"size": 11, "color": DEEP}),
        ("3. Не определён чёткий сценарий — что именно двойник должен оптимизировать",
         {"size": 11, "color": DEEP}),
    ], line_spacing=1.3)
    add_notes(slide, "Southeast Asian Port digital twin: 12 миллионов долларов, 18 месяцев, списан в 2024 году. Кейс задокументирован в отраслевом отчёте Build in Digital [41] — название порта анонимизировано из соображений конфиденциальности. Три корневые причины. Первая — фрагментированные данные из 5 источников без единой схемы: датчики на кранах, GPS контейнеров, ERP-система, погодные сервисы, системы манифестов — все разные форматы, разные частоты. Вторая — чрезмерный акцент на 3D-визуализации: бюджет ушёл на красивую сцену в Unity, а не на конвейер данных. Третья — не был определён чёткий сценарий применения: проект назывался «двойник для порта», но никто не мог сказать, какое конкретное решение двойник должен принимать. Урок: аудит слоя данных обязателен ПЕРЕД запуском пилота. Если хотя бы один ответ «нет» в чек-листе из 5 вопросов — деньги на ветер.")


def s10_data_audit(p):
    """s10 — 5-question data layer audit checklist."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Аудит слоя данных: 5 вопросов до пилота",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Если хотя бы один ответ «нет» — пилот = деньги на ветер (урок анонимного порта 2024)",
             size=13, italic=True, color=LIGHT)
    questions = [
        ("1", "Доступ к историческим данным ≥1 год?",
         "Без этого машинное обучение не отличает сезонные сдвиги от долгосрочных трендов"),
        ("2", "Частота опроса датчиков ≥10× выше полосы управления?",
         "По теореме Найквиста; для регулирования с шагом 1 мин нужны данные с 6-сек интервалом"),
        ("3", "Происхождение разметки документировано?",
         "Кто, когда, по какому критерию ставил метку «дефект» / «норма» / «отказ»"),
        ("4", "Дрейф датчиков калибруется и журналируется?",
         "Без калибровки модель учится на дрейфующих данных и быстро деградирует"),
        ("5", "Назначен ответственный за управление данными?",
         "Кто отвечает за схему, доступ, срок хранения, конфиденциальность"),
    ]
    y_start = 1.8
    item_h = 0.95
    gap = 0.08
    for i, (num, q, desc) in enumerate(questions):
        y = y_start + i * (item_h + gap)
        rounded_box(slide, 0.5, y, 12.33, item_h)
        # Number circle
        circle(slide, 0.7, y + 0.15, 0.65, 0.65, fill=GOLD)
        text_box(slide, 0.7, y + 0.15, 0.65, 0.65, num,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Question
        text_box(slide, 1.55, y + 0.1, 10.6, 0.45, q,
                 size=15, bold=True, color=DEEP)
        # Desc
        text_box(slide, 1.55, y + 0.55, 10.6, 0.4, desc,
                 size=11, italic=True, color=DARK_GREY)
    footer(slide, "Аудит слоя данных — обязательная проверка перед запуском любого пилота")


# ====================================================================
# SECTION 2 — A0 (§2) — s11 divider + s12-s14
# ====================================================================

def s11_div(p):
    section_divider(p, 2, 2, "A0 — Наблюдать",
                    "Vision QC · прогностическое обслуживание · границы применимости",
                    "Раздел 2 · нижняя ступень шкалы автономии")


def s12_vision(p):
    """s12 — Vision QC + FP cascade."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.35, 12.33, 0.7,
             "Машинное зрение: ≥99% при 0,1–2% ложных срабатываний",
             size=26, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.05, 12.33, 0.55,
             "Но даже 1% ложных срабатываний × 10 000 деталей = 100 годных отвергнуто за смену — каскад срабатываний",
             size=13, italic=True, color=LIGHT)
    # Left: chart
    chart = ASSETS / "charts" / "s12-fp-cascade.png"
    if chart.exists():
        rounded_box(slide, 0.5, 1.7, 7.5, 5.0)
        add_image(slide, chart, 0.7, 1.85, 7.1, 4.7)
    # Right: cost-of-FP card
    rounded_box(slide, 8.2, 1.7, 4.6, 5.0, fill=GOLD_TINT, stroke=GOLD, stroke_w=2.0)
    multiline(slide, 8.4, 1.85, 4.2, 4.8, [
        ("100", {"size": 72, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("годных деталей отвергнуто",
         {"size": 13, "bold": True, "color": DEEP, "align": PP_ALIGN.CENTER}),
        ("за смену: 1% × 10 000",
         {"size": 12, "italic": True, "color": DARK_GREY, "align": PP_ALIGN.CENTER}),
        ("", {"size": 12}),
        ("→ ручная переборка", {"size": 12, "color": DEEP}),
        ("→ растёт стоимость пересортировки", {"size": 12, "color": DEEP}),
        ("→ падает пропускная способность", {"size": 12, "color": DEEP}),
        ("→ оператор начинает игнорировать AI", {"size": 12, "color": DEEP}),
        ("→ доверие к AI рушится", {"size": 12, "bold": True, "color": GOLD}),
    ], line_spacing=1.35)
    attribution(slide, "Indus Vision / Jidoka 2026 · Overview.ai 2026 · бенчмарки вендоров для отлаженных систем")
    add_notes(slide, "Запомните арифметику. Vision QC в отлаженных системах достигает 99% и выше точности при ложных срабатываниях от 0,1 до 2%. Это лучший класс для дискретного производства. Но даже 1% ложных срабатываний при 10 000 деталях за смену — это 100 годных отвергнуто. Что дальше: ручная переборка, растёт стоимость пересортировки, падает пропускная способность, оператор начинает игнорировать сигналы AI, доверие к AI рушится. Это не «редкая проблема» — это структурный риск Vision QC. Альтернатива для нестабильного процесса: стабилизация процесса ПЕРЕД Vision AI. Для жёстких допусков ±0,001 мм: метрология + GD&T + SPC.")


def s13_pdm(p):
    """s13 — PdM ROI breakdown."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Прогностическое обслуживание: ROI 10:1 за 2 года",
             size=28, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Программа $200K–$600K → $1,2M–$3,5M годовой экономии · окупаемость 18–36 месяцев",
             size=13, italic=True, color=LIGHT)
    chart = ASSETS / "charts" / "s13-pdm-effects.png"
    if chart.exists():
        rounded_box(slide, 0.5, 1.7, 7.5, 5.0)
        add_image(slide, chart, 0.7, 1.85, 7.1, 4.7)
    # Right: case-anchors
    rounded_box(slide, 8.2, 1.7, 4.6, 5.0, fill=SURFACE, stroke=GOLD, stroke_w=1.5)
    multiline(slide, 8.4, 1.85, 4.2, 4.8, [
        ("ЯКОРНЫЕ КЕЙСЫ:", {"size": 11, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}),
        ("", {"size": 8}),
        ("Цементный завод", {"size": 14, "bold": True, "color": DEEP}),
        ("57× окупаемость за 6 месяцев",
         {"size": 16, "bold": True, "color": GOLD}),
        ("мониторинг без замены оборудования", {"size": 10, "italic": True, "color": SLATE}),
        ("", {"size": 8}),
        ("Химический завод", {"size": 14, "bold": True, "color": DEEP}),
        ("$2 миллиона", {"size": 16, "bold": True, "color": MID}),
        ("ежегодной экономии", {"size": 11, "italic": True, "color": DARK_GREY}),
        ("", {"size": 8}),
        ("Программа прогн. обслуживания:", {"size": 11, "bold": True, "color": MID}),
        ("$200K–$600K инвестиций →", {"size": 11, "color": DEEP}),
        ("$1,2M–$3,5M экономии в год →", {"size": 11, "color": DEEP}),
        ("окупаемость 18–36 месяцев", {"size": 11, "bold": True, "color": GOLD}),
    ], line_spacing=1.3)
    attribution(slide, "Deloitte 2026 / oxmaint 2026 · средние значения по индустрии")


def s14_a0_limits(p):
    """s14 — When vision/PdM NOT applicable."""
    slide = blank(p)
    set_bg(slide, WHITE)
    text_box(slide, 0.5, 0.4, 12.33, 0.7,
             "Где A0 не применим: границы зрения и прогн. обслуживания",
             size=24, bold=True, color=DEEP)
    text_box(slide, 0.5, 1.15, 12.33, 0.4,
             "Жёсткие допуски → метрология + GD&T + SPC · редкие отказы → физика + RCM",
             size=13, italic=True, color=LIGHT)
    # Two columns
    col_y = 1.8
    col_h = 4.8
    col_w = 6.0
    # Left: Tight tolerances
    rounded_box(slide, 0.5, col_y, col_w, col_h, fill=SURFACE, stroke=MID, stroke_w=2.0)
    rectangle(slide, 0.5, col_y, col_w, 0.7, fill=MID)
    text_box(slide, 0.5, col_y + 0.08, col_w, 0.55, "Жёсткие допуски ±0,001 мм",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    multiline(slide, 0.7, col_y + 0.9, col_w - 0.4, 3.8, [
        ("Что происходит:", {"size": 12, "bold": True, "color": LIGHT}),
        ("Камера разрешает 0,01–0,1 мм; не видит микрометровых отклонений.",
         {"size": 12, "color": DEEP}),
        ("", {"size": 8}),
        ("АЛЬТЕРНАТИВА:", {"size": 13, "bold": True, "color": GOLD}),
        ("· Метрология (лазерный сканер, КИМ)", {"size": 12, "color": DEEP}),
        ("· GD&T — геометрические допуски",
         {"size": 12, "color": DEEP}),
        ("· SPC — статистическое управление процессом",
         {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Якорь: «AI не дотягивает физическое разрешение».",
         {"size": 11, "italic": True, "color": MID}),
    ], line_spacing=1.35)
    # Right: Rare events
    rounded_box(slide, 6.83, col_y, col_w, col_h, fill=SURFACE, stroke=TEAL, stroke_w=2.0)
    rectangle(slide, 6.83, col_y, col_w, 0.7, fill=TEAL)
    text_box(slide, 6.83, col_y + 0.08, col_w, 0.55, "Редкие отказы MTBF >1 года",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    multiline(slide, 7.03, col_y + 0.9, col_w - 0.4, 3.8, [
        ("Что происходит:", {"size": 12, "bold": True, "color": LIGHT}),
        ("Выборка отказов <30 — машинное обучение без статистики не учится; галлюцинирует закономерности в шуме.",
         {"size": 12, "color": DEEP}),
        ("", {"size": 8}),
        ("АЛЬТЕРНАТИВА:", {"size": 13, "bold": True, "color": GOLD}),
        ("· Физическая симуляция (МКЭ, износ)", {"size": 12, "color": DEEP}),
        ("· RCM — обслуживание, ориентированное на надёжность",
         {"size": 12, "color": DEEP}),
        ("· Регламент по физике, не календарю", {"size": 12, "color": DEEP}),
        ("", {"size": 6}),
        ("Якорь: «методология Nowlan–Heap 1978 из авиации работает там, где ML — нет».",
         {"size": 11, "italic": True, "color": MID}),
    ], line_spacing=1.35)
    # Bottom takeaway
    rounded_box(slide, 0.5, 6.75, 12.33, 0.4, fill=GOLD_TINT, stroke=GOLD, stroke_w=1.5)
    text_box(slide, 0.7, 6.78, 12.0, 0.35,
             "MTBF (среднее время между отказами) · GD&T · SPC · RCM — четыре альтернативных инструмента вместо AI",
             size=11, italic=True, color=DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


print(f"Loaded part 1 of build_lec12.py")
